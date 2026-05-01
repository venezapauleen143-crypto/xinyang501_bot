"""
Claude Code 工具腳本 - 整合小牛馬所有技能
用法：python claude_tools.py <tool> [參數...]

tools:
  weather <城市>                查詢天氣
  stock <代號> [期間]           查詢股票（期間: 1d 5d 1mo 3mo 6mo 1y，預設 1mo）
  image <prompt> [文字]         生成圖片（prompt 用英文，文字為疊加中文）
  screenshot                    截圖存到桌面
  click <x> <y>                 點擊
  double_click <x> <y>          雙擊
  right_click <x> <y>           右鍵點擊
  move <x> <y>                  移動滑鼠
  type <文字>                   輸入文字
  press <按鍵>                  按下按鍵
  open <程式>                   開啟程式
  scroll <up|down> [格數]       滾動
  pos                           取得目前滑鼠座標
  bot_status                    查看 bot 是否在執行
  bot_restart                   重啟 bot
  schedule_list                 列出所有排程任務
  schedule_add <名稱> <時間HH:MM> <腳本路徑>   新增每日排程（並設定可喚醒）
  schedule_del <名稱>           刪除排程任務
  memory_save <chat_id> <內容>  儲存長期記憶
  memory_list <chat_id>         列出長期記憶
  memory_del <chat_id> <id>     刪除指定記憶
  vision [問題]                 截圖並讓 Claude 分析畫面內容
  find_image <圖片路徑>         在螢幕上找到指定圖片的位置
  browser <動作> [參數]         瀏覽器自動化（open/click/type/get_text/screenshot/goto/close）
  window_list                   列出所有視窗
  window_focus <標題關鍵字>     切換到指定視窗
  window_close <標題關鍵字>     關閉指定視窗
  window_min <標題關鍵字>       最小化視窗
  window_max <標題關鍵字>       最大化視窗
  hotkey <按鍵組合>             執行組合鍵（如 ctrl+c, alt+tab）
  clipboard_get                 讀取剪貼簿內容
  clipboard_set <內容>          寫入剪貼簿
  file_list <路徑>              列出資料夾內容
  file_read <路徑>              讀取文字檔案
  file_write <路徑> <內容>      寫入文字檔案
  file_delete <路徑>            刪除檔案或資料夾
  file_copy <來源> <目標>       複製檔案
  file_move <來源> <目標>       移動/重新命名檔案
  file_search <資料夾> <關鍵字> 搜尋檔案
  sysinfo                       查看 CPU/記憶體/磁碟使用率
  process_list                  列出所有執行中程序
  process_kill <名稱或PID>      結束程序
  notify <標題> <訊息>          發送 Windows 桌面通知
  tts <文字>                    文字轉語音朗讀
  record_start                  開始錄製滑鼠鍵盤動作
  record_stop                   停止錄製並儲存
  record_play <檔案>            重播錄製的動作
  email <收件人> <主旨> <內容>  發送 Email
  stt                           語音辨識（麥克風錄音轉文字）
  ocr [圖片路徑]                OCR 截圖或指定圖片，辨識文字
  workflow_run <json檔>         執行自動化工作流程
  workflow_save <名稱> <json>   儲存工作流程
  screen_watch <圖片> <指令>    監控螢幕出現指定圖片時執行指令
  monitors                      列出所有螢幕資訊
  zip <來源> <目標zip>          壓縮檔案或資料夾
  unzip <zip檔> <目標資料夾>    解壓縮
  download <網址> [儲存路徑]    下載檔案
  print_file <檔案路徑>         列印文件
  wifi_list                     列出可用 WiFi 網路
  wifi_connect <SSID> <密碼>    連線 WiFi
  screen_stream <秒數>          截圖串流（每秒截圖存桌面，持續N秒）
  wake_listen <關鍵字>          等待語音說出關鍵字後執行
  drag <x1> <y1> <x2> <y2>     拖曳從(x1,y1)到(x2,y2)
  right_menu <x> <y> <項目>     右鍵選單並選擇項目
  ai_plan <目標>                AI 自動規劃並執行多步驟任務
  clipboard_history             顯示剪貼簿歷史紀錄
  vdesktop <left|right|new>     切換虛擬桌面
  power <sleep|restart|shutdown> 電源管理
  bt_scan                       掃描藍牙裝置
  bt_connect <MAC位址>          連線藍牙裝置
  run_python <程式碼>           直接執行 Python 程式碼
  run_shell <指令>              執行 PowerShell 指令
  word_read <路徑>              讀取 Word 文件
  word_write <路徑> <內容>      寫入 Word 文件
  excel_read <路徑> [工作表]    讀取 Excel
  excel_write <路徑> <工作表> <資料JSON> 寫入 Excel
  pdf_read <路徑>               讀取 PDF 文字
  screen_diff <間隔秒> <區域>   偵測螢幕區域變化
  scrape <網址> <CSS選擇器>     爬取網頁指定內容
  img_edit <動作> <路徑> [參數] 圖片編輯（crop/resize/text/merge）
  gdrive_upload <檔案> <資料夾ID> 上傳到 Google Drive
  gdrive_download <檔案ID> <路徑> 從 Google Drive 下載
  db_query <資料庫路徑> <SQL>   執行 SQLite 查詢
  db_mysql <host> <db> <SQL>    執行 MySQL 查詢
  encrypt <檔案路徑> <密碼>     加密檔案
  decrypt <檔案路徑> <密碼>     解密檔案
  clipboard_watch <秒數>        監控剪貼簿變化
  qr_gen <內容> [路徑]          生成 QR Code
  qr_scan [圖片路徑]            掃描 QR Code（截圖或指定圖片）
  screen_record [秒數] [輸出路徑]  螢幕錄影存 mp4
  webcam [輸出路徑]             攝影機拍照
  translate <文字> [目標語言] [來源語言]  翻譯（目標預設 zh-TW）
  chart <類型> <資料JSON> [標題] [輸出路徑]  生成圖表（line/bar/pie）
  pptx_read <路徑>              讀取 PowerPoint
  pptx_create <路徑> <slides_json>  建立 PowerPoint
  api_call <方法> <URL> [headers_json] [body_json]  呼叫 REST API
  watchdog <程序名> <腳本路徑> [秒數]  守護程序，崩潰自動重啟
  ssh_run <host> <user> <password> <指令>  SSH 執行遠端指令
  sftp_upload <host> <user> <pass> <local> <remote>  SFTP 上傳
  sftp_download <host> <user> <pass> <remote> <local>  SFTP 下載
  net_ping <host> [次數]        Ping 網路主機
  net_traceroute <host>         路由追蹤
  net_portscan <host> [ports]   Port 掃描
  win_service <list|start|stop> [服務名]  Windows 服務管理
  pdf_merge <paths_json> <輸出路徑>  合併 PDF
  pdf_split <路徑> <輸出資料夾>  分割 PDF 每頁
  pdf_watermark <路徑> <文字> [輸出路徑]  PDF 加浮水印
  audio_convert <輸入> <輸出>   音訊格式轉換
  audio_trim <輸入> <起始ms> <結束ms> [輸出]  音訊剪輯
  discord_notify <webhook_url> <訊息>  Discord 推播
  line_notify <token> <訊息>    LINE Notify 推播
  disk_clean <list|clean>       磁碟暫存清理
  backup <來源> <目標資料夾>    備份壓縮
  registry_read <key_path> [value_name]  讀取登錄檔
  registry_write <key_path> <value_name> <value>  寫入登錄檔
  video_screenshot <路徑> [秒數] [輸出]  影片截取畫面
  video_trim <路徑> <起始秒> <結束秒> [輸出]  影片剪輯
  monitor_list                  列出所有螢幕資訊
  email_read <host> <user> <pass> [資料夾] [數量]  讀取 IMAP 郵件
  gcal_list [天數]              列出 Google Calendar 行程
  gcal_add <標題> <開始> <結束> [說明]  新增行程（時間格式 2026-04-13T10:00:00）
  global_hotkey <熱鍵> <指令> [秒數]  監聽全域快捷鍵
  git_op <動作> [repo] [參數]   Git 操作（status/log/pull/add/commit/push/diff）
  hw_monitor                    硬體監控（GPU/電池/溫度）
  report_gen <標題> <資料JSON> [輸出路徑]  生成 HTML 報告
  dropbox_upload <local> <remote> [token]  上傳到 Dropbox
  dropbox_download <remote> <local> [token]  從 Dropbox 下載
  docker_op <動作> [容器名]     Docker 操作（list/start/stop/logs/images）
  pdf_to_images <路徑> [輸出資料夾] [dpi]  PDF 轉圖片
  barcode_scan [圖片路徑]       掃描條碼或 QR Code
  nlp_summarize <文字>          AI 文字摘要
  nlp_sentiment <文字>          AI 情緒分析
  vpn_control <list|connect|disconnect> [名稱] [帳號] [密碼]  VPN 控制
  sys_restore <create|list> [說明]  Windows 系統還原點
  disk_analyze [路徑] [數量]    磁碟空間分析
  face_detect [圖片路徑] [輸出路徑]  人臉偵測
  video_to_gif <路徑> [起始秒] [持續秒] [輸出] [fps]  影片轉 GIF
  excel_chart <路徑> <工作表> [類型] [標題]  Excel 生成圖表
  speedtest                     網路速度測試
  screenshot_compare [圖1] [圖2] [輸出]  截圖比對差異
  set_reminder <HH:MM或秒數> <訊息>  設定一次性提醒
  webpage_screenshot <網址> [輸出]  網頁全頁截圖
  web_monitor <網址> [selector] [間隔秒] [持續秒]  網頁變化監控
  batch_rename <資料夾> <正規表達式> <替換> [副檔名]  批次重新命名
  img_compress <路徑> [品質0-100] [輸出]  壓縮圖片
  batch_img_process <資料夾> <resize|compress> [寬] [高] [品質]  批次圖片處理
  ocr_translate [圖片路徑] [目標語言]  OCR辨識+翻譯
  ip_info [IP地址]              查詢 IP 地理位置（留空查自己）
  currency <金額> <來源幣> <目標幣>  匯率換算
  event_log [日誌名] [等級] [數量]  Windows 事件日誌
  tts_edge <文字> [語音] [輸出]  微軟 Edge TTS 語音合成
  tts_voices                    列出可用 Edge TTS 語音
  send_email_attach <收件人> <主旨> <內容> [附件路徑,...]  含附件郵件
  clipboard_img_get [輸出路徑]  讀取剪貼簿圖片
  clipboard_img_set <圖片路徑>  複製圖片到剪貼簿
  usb_list                      列出 USB 裝置
  firewall <list|add|remove> [名稱] [方向] [port] [協定]  防火牆管理
  todo <add|list|done|delete|clear> [任務] [id]  任務清單
  file_sync <來源> <目標> [dry_run]  資料夾同步
  sysres_chart [秒數] [輸出]    系統資源使用圖表
  password_save <網站> <帳號> <密碼> <主密碼>  加密儲存密碼
  password_get <網站> <主密碼>  查詢已儲存密碼
  rdp_connect <host> [user] [寬] [高]  RDP 遠端桌面
  chrome_bookmarks              列出 Chrome 書籤
  printer_list                  列出印表機
  printer_jobs                  列出列印佇列
  net_share <list|connect|disconnect> [共享路徑] [磁碟代號] [帳號] [密碼]  網路芳鄰
  font_list [關鍵字]            列出系統字型
  volume_get                    查詢系統音量
  volume_set <0-100>            設定系統音量
  volume_mute [true|false]      靜音/取消靜音
  brightness_get                查詢螢幕亮度
  brightness_set <0-100>        設定螢幕亮度
  resolution_list               查詢螢幕解析度
  media_control <動作>          媒體控制（play_pause/next/prev/stop/volume_up/volume_down/mute）
  audio_devices                 列出音訊裝置
  audio_switch <裝置名稱>       切換音訊輸出裝置
  software_list [關鍵字]        列出已安裝軟體
  software_install <名稱>       安裝軟體（winget）
  software_uninstall <名稱>     卸載軟體
  startup_list                  列出開機自啟動程式
  startup_add <名稱> <指令>     新增開機自啟動
  startup_remove <名稱>         移除開機自啟動
  env_get [變數名]              查詢環境變數
  env_set <名稱> <值> [true]    設定環境變數（true=永久）
  user_list                     列出使用者帳戶
  user_create <帳號> <密碼>     建立使用者
  user_delete <帳號>            刪除使用者
  win_update <list|install|check>  Windows 更新管理
  device_list [關鍵字]          裝置管理員列表
  device_toggle <名稱> [true|false]  啟用/停用裝置
  netadapter_list               列出網路介面卡
  netadapter_toggle <名稱> [true|false]  啟用/停用網路卡
  dns_get                       查詢 DNS 設定
  dns_set <介面> <DNS1> [DNS2]  設定 DNS
  ip_config <介面> <IP> [遮罩] [閘道]  設定靜態 IP
  hosts_list                    列出 hosts 設定
  hosts_add <IP> <domain>       新增 hosts
  hosts_remove <domain>         移除 hosts
  net_traffic [秒數]            監控網路流量
  if_then <條件類型> <條件值> <指令> [秒數]  條件式自動化
  window_arrange <side_by_side|quad|stack|maximize_all>  多視窗排列
  region_ocr <x> <y> <w> <h> [語言]  指定區域 OCR
  window_screenshot <視窗標題關鍵字> [輸出路徑]  指定視窗截圖

  ── 新增工具（同步 bot.py）──
  analyze_pdf <路徑> [最大字數]    分析 PDF 文件
  audio_process <action> <輸入> [輸出] [起始ms] [結束ms]  音訊處理
  auto_skill <action> [目標] [名稱] [code]  自動技能管理
  auto_trade <action> [代號] [數量] [價格] [類型]  自動交易
  automation <action> [參數...]    自動化操作
  barcode [圖片路徑]               掃描條碼
  bluetooth <action> [MAC]         藍牙操作
  browser_advanced <action> [參數]  進階瀏覽器控制
  browser_control <action> [url] [selector] [text]  瀏覽器控制
  calendar <action> [天數] [標題] [開始] [結束]  Google 日曆
  clipboard <action> [文字]        剪貼簿操作
  clipboard_image <action> [路徑]  剪貼簿圖片
  cloud_storage <action> <路徑> [drive_id]  雲端儲存
  compare_stocks <代號,...> [指標]  比較股票
  database <type> <db> <SQL>       資料庫操作
  ddg_search <關鍵字> [地區] [數量]  DuckDuckGo 搜尋
  desktop_control <action> [x] [y] [text] [app]  桌面控制
  device_manager <action> [名稱]   裝置管理員
  disk_backup <action> [來源] [目標]  磁碟備份
  display <action> [亮度]          顯示設定
  docker <action> [容器名]         Docker 操作
  document_control <action> <路徑> [內容]  文件控制
  download_file <URL> [路徑]       下載檔案
  dropbox <action> <本地> <遠端> [token]  Dropbox
  email_control <host> <user> <pass>  Email 控制
  emotion_detect <action> [文字] [圖片]  情緒偵測
  encrypt_file <action> <路徑> <密碼>  加解密檔案
  env_var <action> [名稱] [值]     環境變數
  file_system <action> [路徑] [目標]  檔案系統操作
  file_tools <action> <路徑>       檔案工具
  file_transfer <action> <來源> [目標]  檔案傳輸
  find_image_on_screen <圖片> [信心值]  螢幕找圖
  generate_image <prompt> [寬] [高] [疊加文字]  生成圖片（含尺寸）
  get_candlestick_chart <代號> [期間]  K線圖分析
  get_crypto <幣種> [vs貨幣]       查詢加密貨幣
  get_earnings <代號>              查詢財報
  get_etf <代號>                   查詢 ETF
  get_finance_news [來源] [數量]   財經新聞
  get_forex <貨幣對>               查詢匯率
  get_fundamentals <代號>          查詢基本面
  get_macro <指標>                 總經指標
  get_market_sentiment             市場情緒
  get_stock <代號> [期間]          查詢股票（進階版）
  get_stock_advanced <代號> [指標]  進階技術分析
  get_weather <城市>               查詢天氣（進階版）
  git <action> [repo] [message]    Git 操作
  goal_manager <action> [目標]     目標管理
  google_trends <關鍵字,...> [時間] [地區]  Google Trends
  hardware                         硬體監控
  image_edit <action> <路徑> [參數]  圖片編輯
  image_tools <action> [路徑] [quality]  圖片工具
  knowledge_base <action> [content] [query]  知識庫
  long_term_memory <action> [chat_id] [內容]  長期記憶
  lookup <ip|currency> [參數]      查詢工具
  manage_schedule <action> [名稱] [時間]  排程管理
  media <action> [裝置名]          媒體控制
  monitor_config                   螢幕設定資訊
  multi_deploy <action> [host] [user]  多機部署
  multi_perspective <主題>         多角度分析
  network_config <action> [名稱]   網路設定
  network_diag <action> <host>     網路診斷
  news_monitor <action> [關鍵字]   新聞監控
  news_search <關鍵字> [語言] [數量]  新聞搜尋
  nlp <action> <文字>              NLP 工具
  osint_search <action> [query]    OSINT 搜尋
  password_mgr <action> <站> <主密碼>  密碼管理
  pdf_edit <action> [路徑] [輸出]  PDF 編輯
  pdf_image <路徑> [輸出] [dpi]    PDF 轉圖片
  pentest <action> [target]        滲透測試
  portfolio <action> [chat_id] [symbol]  投資組合
  power_control <action>           電源控制
  pptx_control <action> <路徑>     PowerPoint 控制
  proactive_alert <action> [名稱]  主動警報
  ptt_search <關鍵字> [看板] [數量]  PTT 搜尋
  push_notify <platform> <訊息> <webhook>  推播通知
  qr_code <action> [content] [path]  QR Code 操作
  read_screen [問題] [螢幕號]      讀取螢幕內容
  read_webpage <URL> [最大字數]    讀取網頁內容
  registry <action> <key> [value_name]  登錄檔操作
  reminder <時間> <訊息>           設定提醒
  report <標題> <資料JSON> [輸出]  報告生成
  restore_point <action> [說明]    系統還原點
  run_code <type> <code>           執行程式碼
  screen_vision [問題]             螢幕視覺分析
  scroll_at [方向] [格數] [x] [y] [螢幕]  指定位置滾動
  self_benchmark <action>          自我評測
  send_email <收件人> <主旨> <內容>  發送 Email（進階版）
  send_voice <文字> [語音]         生成語音檔
  smart_home <action> [device]     智慧家居
  software <action> [名稱]         軟體管理
  ssh_sftp <action> <host> <user> <pass>  SSH/SFTP
  startup <action> [名稱] [指令]   開機自啟
  system_monitor <action> [target]  系統監控
  system_tools <action>            系統工具集
  think_as <人物> <問題>           角色思考
  threat_intel <action> [target]   威脅情報
  todo_list <action> [task] [id]   任務清單
  tts_advanced <action> [文字] [語音]  進階 TTS
  user_account <action> [帳號]     使用者帳戶
  video_gif <路徑> [起始秒] [持續秒]  影片轉 GIF
  video_process <action> <路徑>    影片處理
  virtual_desktop <action>         虛擬桌面
  voice_cmd <action> [持續秒]      語音命令
  voice_id <action> [名稱] [音檔]  聲紋辨識
  volume <action> [音量]           音量控制
  vpn <action> [名稱]              VPN 控制
  wait_seconds <秒數>              等待指定秒數
  web_scrape <action> [url] [selector]  網頁爬取
  webpage_shot <action> <url>      網頁截圖
  wikipedia_search <關鍵字> [語言]  Wikipedia 搜尋
  win_notify_relay <action>        Windows 通知轉發
  window_control <action> [關鍵字]  視窗控制
  window_manager <action> [視窗名]  視窗管理員
  windows_update <action>          Windows 更新管理
  workflow <action> [名稱] [steps]  工作流程
  youtube_summary <URL>            YouTube 字幕摘要
"""

import sys
import os
import io
import time
import sqlite3
import subprocess
import requests
import pyautogui
from pathlib import Path
from dotenv import load_dotenv

# 強制 stdout/stderr 使用 UTF-8（修 subprocess 跑時 emoji GBK encode 錯誤）
# reconfigure 在 console 跑 work，但 subprocess pipe 不一定 work，所以用 TextIOWrapper 強制 wrap
try:
    if hasattr(sys.stdout, 'buffer'):
        sys.stdout = io.TextIOWrapper(sys.stdout.buffer, encoding='utf-8', errors='replace', line_buffering=True)
    if hasattr(sys.stderr, 'buffer'):
        sys.stderr = io.TextIOWrapper(sys.stderr.buffer, encoding='utf-8', errors='replace', line_buffering=True)
except Exception:
    pass

load_dotenv(Path("C:/Users/blue_/claude-telegram-bot/.env"))

pyautogui.FAILSAFE = True

# ── Additional imports for extracted bot functions ──
import json
import logging
import datetime as dt
import threading
import urllib.parse
import traceback

SCREENSHOT_DIR = Path.home() / "Desktop" / "測試檔案"

# ── 天氣 ────────────────────────────────────────────

def get_weather(city: str):
    try:
        res = requests.get(
            f"https://wttr.in/{city}?format=j1",
            headers={"User-Agent": "curl/7.68.0"},
            timeout=10
        )
        data = res.json()
        current = data["current_condition"][0]
        area = data["nearest_area"][0]
        location = area["areaName"][0]["value"] + ", " + area["country"][0]["value"]
        desc = current["lang_zh-tw"][0]["value"] if "lang_zh-tw" in current else current["weatherDesc"][0]["value"]
        temp = current["temp_C"]
        feels = current["FeelsLikeC"]
        humidity = current["humidity"]
        wind = current["windspeedKmph"]
        wind_dir = current["winddir16Point"]
        print(
            f"📍 {location}\n"
            f"🌤 {desc}\n"
            f"🌡 氣溫：{temp}°C（體感 {feels}°C）\n"
            f"💧 濕度：{humidity}%\n"
            f"💨 風速：{wind} km/h（{wind_dir}）"
        )
    except Exception as e:
        print(f"查不到「{city}」的天氣：{e}")


# ── 股票 ────────────────────────────────────────────

def calc_rsi(closes, period=14):
    deltas = closes.diff().dropna()
    gains = deltas.clip(lower=0)
    losses = -deltas.clip(upper=0)
    avg_gain = gains.rolling(period).mean().iloc[-1]
    avg_loss = losses.rolling(period).mean().iloc[-1]
    if avg_loss == 0:
        return 100.0
    rs = avg_gain / avg_loss
    return round(100 - (100 / (1 + rs)), 1)

def get_stock(symbol: str, period: str = "1mo"):
    try:
        import yfinance as yf
        ticker = yf.Ticker(symbol)
        info = ticker.info
        hist = ticker.history(period="3mo")

        if hist.empty:
            print(f"找不到「{symbol}」的股票數據，請確認代號是否正確。")
            return

        name = info.get("longName") or info.get("shortName") or symbol
        currency = info.get("currency", "")
        current = hist["Close"].iloc[-1]
        prev = hist["Close"].iloc[-2] if len(hist) > 1 else current
        change = current - prev
        change_pct = (change / prev) * 100 if prev else 0
        arrow = "▲" if change >= 0 else "▼"
        volume = hist["Volume"].iloc[-1]
        avg_volume = hist["Volume"].tail(20).mean()
        volume_ratio = volume / avg_volume if avg_volume else 1

        ma5 = hist["Close"].tail(5).mean()
        ma20 = hist["Close"].tail(20).mean()
        ma60 = hist["Close"].tail(60).mean() if len(hist) >= 60 else hist["Close"].mean()
        rsi = calc_rsi(hist["Close"]) if len(hist) >= 15 else None

        if ma5 > ma20 > ma60:
            trend = "強勢多頭（MA5>MA20>MA60）📈"
        elif ma5 < ma20 < ma60:
            trend = "強勢空頭（MA5<MA20<MA60）📉"
        elif ma5 > ma20:
            trend = "短線偏多（MA5>MA20）🔼"
        else:
            trend = "短線偏空（MA5<MA20）🔽"

        rsi_note = ""
        if rsi is not None:
            if rsi >= 80:
                rsi_note = "（嚴重超買 ⚠️）"
            elif rsi >= 70:
                rsi_note = "（超買區間）"
            elif rsi <= 20:
                rsi_note = "（嚴重超賣 💡）"
            elif rsi <= 30:
                rsi_note = "（超賣區間）"
            else:
                rsi_note = "（中性）"

        market_cap = info.get("marketCap")
        pe_ratio = info.get("trailingPE")
        week52_high = info.get("fiftyTwoWeekHigh")
        week52_low = info.get("fiftyTwoWeekLow")
        high_period = hist["High"].tail(20).max()
        low_period = hist["Low"].tail(20).min()
        pct_from_high = ((current - high_period) / high_period * 100) if high_period else 0
        pct_from_low = ((current - low_period) / low_period * 100) if low_period else 0

        result = (
            f"📊 {name} ({symbol})\n"
            f"💰 現價：{current:.2f} {currency}  {arrow} {abs(change):.2f} ({change_pct:+.2f}%)\n"
            f"📦 成交量：{volume:,}（均量 {volume_ratio:.1f}x）\n"
            f"\n── 技術指標 ──\n"
            f"MA5：{ma5:.2f}　MA20：{ma20:.2f}　MA60：{ma60:.2f}\n"
            f"趨勢：{trend}\n"
        )
        if rsi is not None:
            result += f"RSI(14)：{rsi}{rsi_note}\n"
        result += (
            f"近20日高點：{high_period:.2f}（距高 {pct_from_high:.1f}%）\n"
            f"近20日低點：{low_period:.2f}（距低 +{pct_from_low:.1f}%）\n"
        )
        if week52_high and week52_low:
            result += f"52週高低：{week52_low:.2f} ~ {week52_high:.2f}\n"
        result += "\n── 基本面 ──\n"
        if market_cap:
            mc_str = f"{market_cap/1e12:.2f}兆" if market_cap >= 1e12 else f"{market_cap/1e8:.0f}億"
            result += f"市值：{mc_str} {currency}\n"
        if pe_ratio:
            result += f"本益比：{pe_ratio:.1f}\n"

        print(result.strip())

    except Exception as e:
        print(f"查詢「{symbol}」失敗：{e}")


# ── 圖片生成 ────────────────────────────────────────

def add_text_overlay(image_bytes: bytes, text: str) -> bytes:
    try:
        from PIL import Image, ImageDraw, ImageFont
        img = Image.open(io.BytesIO(image_bytes)).convert("RGBA")
        w, h = img.size
        overlay = Image.new("RGBA", img.size, (0, 0, 0, 0))
        draw = ImageDraw.Draw(overlay)
        font_size = max(36, w // 12)
        font_path = "C:/Windows/Fonts/msjhbd.ttc"
        try:
            font = ImageFont.truetype(font_path, font_size)
        except Exception:
            font = ImageFont.load_default()
        bbox = draw.textbbox((0, 0), text, font=font)
        tw, th = bbox[2] - bbox[0], bbox[3] - bbox[1]
        x = (w - tw) // 2
        y = h - th - int(h * 0.08)
        for dx, dy in [(-2,-2),(2,-2),(-2,2),(2,2)]:
            draw.text((x+dx, y+dy), text, font=font, fill=(0, 0, 0, 180))
        draw.text((x, y), text, font=font, fill=(255, 255, 255, 255))
        result = Image.alpha_composite(img, overlay).convert("RGB")
        buf = io.BytesIO()
        result.save(buf, format="JPEG", quality=92)
        return buf.getvalue()
    except Exception:
        return image_bytes

def generate_image(prompt: str, overlay_text: str = ""):
    hf_token = os.getenv("HF_TOKEN")
    if not hf_token:
        print("未設定 HF_TOKEN，請在 .env 加入 HF_TOKEN=your_token")
        return
    headers = {"Authorization": f"Bearer {hf_token}"}
    payload = {"inputs": prompt}
    models = [
        "black-forest-labs/FLUX.1-schnell",
        "stabilityai/stable-diffusion-xl-base-1.0",
    ]
    image_bytes = None
    for model in models:
        for attempt in range(2):
            try:
                print(f"嘗試模型：{model}...")
                res = requests.post(
                    f"https://router.huggingface.co/hf-inference/models/{model}",
                    headers=headers,
                    json=payload,
                    timeout=120
                )
                if res.status_code == 200 and res.headers.get("content-type", "").startswith("image"):
                    image_bytes = res.content
                    break
                if res.status_code == 503:
                    time.sleep(10)
                    continue
            except Exception:
                pass
        if image_bytes:
            break

    if not image_bytes:
        print("圖片生成失敗")
        return

    if overlay_text:
        image_bytes = add_text_overlay(image_bytes, overlay_text)

    filename = SCREENSHOT_DIR / f"generated_{dt.datetime.now().strftime('%Y%m%d_%H%M%S')}.jpg"
    with open(filename, "wb") as f:
        f.write(image_bytes)
    print(f"圖片已儲存：{filename}")


# ── 桌面控制 ────────────────────────────────────────

def screenshot():
    img = pyautogui.screenshot()
    filename = SCREENSHOT_DIR / f"screenshot_{dt.datetime.now().strftime('%Y%m%d_%H%M%S')}.png"
    img.save(filename)
    print(f"截圖已儲存：{filename}")
    return str(filename)

def click(x, y):
    pyautogui.click(int(x), int(y))
    print(f"已點擊 ({x}, {y})")

def double_click(x, y):
    pyautogui.doubleClick(int(x), int(y))
    print(f"已雙擊 ({x}, {y})")

def right_click(x, y):
    pyautogui.rightClick(int(x), int(y))
    print(f"已右鍵點擊 ({x}, {y})")

def move(x, y):
    pyautogui.moveTo(int(x), int(y), duration=0.3)
    print(f"滑鼠已移到 ({x}, {y})")

def type_text(text):
    pyautogui.write(text, interval=0.05)
    print(f"已輸入：{text}")

def press_key(key):
    pyautogui.press(key)
    print(f"已按下：{key}")

def open_app(app):
    subprocess.Popen(app, shell=True)
    print(f"已開啟：{app}")

def scroll(direction, amount=3):
    amount = int(amount)
    pyautogui.scroll(amount if direction == "up" else -amount)
    print(f"已向{direction}滾動 {amount} 格")

def pos():
    x, y = pyautogui.position()
    print(f"目前滑鼠位置：({x}, {y})")


# ── 長期記憶 ────────────────────────────────────────

MEMORY_DB = Path("C:/Users/blue_/claude-telegram-bot/memory.db")

def memory_save(chat_id: int, content: str):
    conn = sqlite3.connect(MEMORY_DB)
    conn.execute("CREATE TABLE IF NOT EXISTS long_term_memory (id INTEGER PRIMARY KEY AUTOINCREMENT, chat_id INTEGER NOT NULL, content TEXT NOT NULL, created_at TIMESTAMP DEFAULT CURRENT_TIMESTAMP)")
    conn.execute("INSERT INTO long_term_memory (chat_id, content) VALUES (?, ?)", (chat_id, content))
    conn.commit()
    conn.close()
    print(f"✅ 已儲存記憶：{content}")

def memory_list(chat_id: int):
    conn = sqlite3.connect(MEMORY_DB)
    conn.execute("CREATE TABLE IF NOT EXISTS long_term_memory (id INTEGER PRIMARY KEY AUTOINCREMENT, chat_id INTEGER NOT NULL, content TEXT NOT NULL, created_at TIMESTAMP DEFAULT CURRENT_TIMESTAMP)")
    rows = conn.execute("SELECT id, content, created_at FROM long_term_memory WHERE chat_id=? ORDER BY id DESC", (chat_id,)).fetchall()
    conn.close()
    if not rows:
        print("目前沒有長期記憶。")
        return
    print(f"📝 長期記憶（chat_id={chat_id}）：")
    for r in rows:
        print(f"  [{r[0]}] {r[1]}  ({r[2]})")

def memory_del(chat_id: int, memory_id: int):
    conn = sqlite3.connect(MEMORY_DB)
    conn.execute("DELETE FROM long_term_memory WHERE id=? AND chat_id=?", (memory_id, chat_id))
    conn.commit()
    conn.close()
    print(f"✅ 已刪除記憶 ID {memory_id}")


# ── Vision 截圖理解 ──────────────────────────────────

def vision(question: str = "請描述這個畫面上有什麼，以及目前電腦在做什麼事。"):
    from anthropic import Anthropic
    import base64
    client = Anthropic()
    img = pyautogui.screenshot()
    buf = io.BytesIO()
    img.save(buf, format="PNG")
    img_b64 = base64.standard_b64encode(buf.getvalue()).decode("utf-8")
    response = client.messages.create(
        model="claude-sonnet-4-6",
        max_tokens=1024,
        messages=[{
            "role": "user",
            "content": [
                {"type": "image", "source": {"type": "base64", "media_type": "image/png", "data": img_b64}},
                {"type": "text", "text": question}
            ]
        }]
    )
    print(response.content[0].text)


# ── 圖像定位 ─────────────────────────────────────────

def find_image(template_path: str, confidence: float = 0.8):
    try:
        location = pyautogui.locateOnScreen(template_path, confidence=confidence)
        if location:
            cx, cy = pyautogui.center(location)
            print(f"✅ 找到圖片：中心座標 ({cx}, {cy})，區域 {location}")
            return cx, cy
        else:
            print("❌ 畫面上找不到該圖片")
            return None
    except Exception as e:
        print(f"❌ 搜尋失敗：{e}")
        return None


# ── 瀏覽器自動化 ──────────────────────────────────────

_browser_context = {}

def browser(action: str, *args):
    from playwright.sync_api import sync_playwright

    if action == "open":
        url = args[0] if args else "https://www.google.com"
        def _open():
            pw = sync_playwright().start()
            b = pw.chromium.launch(headless=False)
            page = b.new_page()
            page.goto(url)
            _browser_context["pw"] = pw
            _browser_context["browser"] = b
            _browser_context["page"] = page
            print(f"✅ 已開啟：{url}")
        _open()

    elif action == "click":
        page = _browser_context.get("page")
        if not page:
            print("❌ 瀏覽器未開啟，請先執行 browser open <url>")
            return
        selector = args[0]
        page.click(selector)
        print(f"✅ 已點擊：{selector}")

    elif action == "type":
        page = _browser_context.get("page")
        if not page:
            print("❌ 瀏覽器未開啟")
            return
        selector, text = args[0], " ".join(args[1:])
        page.fill(selector, text)
        print(f"✅ 已輸入到 {selector}：{text}")

    elif action == "get_text":
        page = _browser_context.get("page")
        if not page:
            print("❌ 瀏覽器未開啟")
            return
        selector = args[0] if args else "body"
        text = page.inner_text(selector)
        print(text[:2000])

    elif action == "screenshot":
        page = _browser_context.get("page")
        if not page:
            print("❌ 瀏覽器未開啟")
            return
        filename = SCREENSHOT_DIR / f"browser_{dt.datetime.now().strftime('%Y%m%d_%H%M%S')}.png"
        page.screenshot(path=str(filename))
        print(f"✅ 截圖已儲存：{filename}")

    elif action == "goto":
        page = _browser_context.get("page")
        if not page:
            print("❌ 瀏覽器未開啟")
            return
        page.goto(args[0])
        print(f"✅ 已前往：{args[0]}")

    elif action == "close":
        if "browser" in _browser_context:
            _browser_context["browser"].close()
            _browser_context["pw"].stop()
            _browser_context.clear()
            print("✅ 瀏覽器已關閉")
    else:
        print(f"❌ 未知動作：{action}。可用：open / click / type / get_text / screenshot / goto / close")


# ── 視窗管理 ─────────────────────────────────────────

def window_list():
    import pygetwindow as gw
    wins = [w for w in gw.getAllWindows() if w.title.strip()]
    for w in wins:
        print(f"  [{w._hWnd}] {w.title}")

def _find_window(keyword):
    import pygetwindow as gw
    wins = [w for w in gw.getAllWindows() if keyword.lower() in w.title.lower()]
    return wins[0] if wins else None

def window_focus(keyword):
    w = _find_window(keyword)
    if w:
        w.activate()
        print(f"✅ 已切換到：{w.title}")
    else:
        print(f"❌ 找不到視窗：{keyword}")

def window_close(keyword):
    w = _find_window(keyword)
    if w:
        w.close()
        print(f"✅ 已關閉：{w.title}")
    else:
        print(f"❌ 找不到視窗：{keyword}")

def window_min(keyword):
    w = _find_window(keyword)
    if w:
        w.minimize()
        print(f"✅ 已最小化：{w.title}")
    else:
        print(f"❌ 找不到視窗：{keyword}")

def window_max(keyword):
    w = _find_window(keyword)
    if w:
        w.maximize()
        print(f"✅ 已最大化：{w.title}")
    else:
        print(f"❌ 找不到視窗：{keyword}")


# ── 組合鍵 / 剪貼簿 ──────────────────────────────────

def hotkey(*keys):
    combo = "+".join(keys)
    pyautogui.hotkey(*keys)
    print(f"✅ 已執行組合鍵：{combo}")

def clipboard_get():
    import pyperclip
    text = pyperclip.paste()
    print(text)

def clipboard_set(text):
    import pyperclip
    pyperclip.copy(text)
    print(f"✅ 已寫入剪貼簿：{text}")


# ── 檔案系統 ─────────────────────────────────────────

def file_list(path="."):
    p = Path(path)
    if not p.exists():
        print(f"❌ 路徑不存在：{path}")
        return
    for item in sorted(p.iterdir()):
        tag = "📁" if item.is_dir() else "📄"
        print(f"  {tag} {item.name}")

def file_read(path):
    p = Path(path)
    if not p.exists():
        print(f"❌ 檔案不存在：{path}")
        return
    print(p.read_text(encoding="utf-8", errors="replace"))

def file_write(path, content):
    Path(path).write_text(content, encoding="utf-8")
    print(f"✅ 已寫入：{path}")

def file_delete(path):
    import shutil
    p = Path(path)
    if not p.exists():
        print(f"❌ 不存在：{path}")
        return
    if p.is_dir():
        shutil.rmtree(p)
    else:
        p.unlink()
    print(f"✅ 已刪除：{path}")

def file_copy(src, dst):
    import shutil
    shutil.copy2(src, dst)
    print(f"✅ 已複製：{src} → {dst}")

def file_move(src, dst):
    import shutil
    shutil.move(src, dst)
    print(f"✅ 已移動：{src} → {dst}")

def file_search(folder, keyword):
    # 如果 keyword 已含 glob 萬用字元 (* ? [) 直接當 pattern 用，否則 wrap 成 *keyword*
    if any(c in keyword for c in "*?["):
        pattern = keyword
    else:
        pattern = f"*{keyword}*"
    try:
        results = list(Path(folder).rglob(pattern))
    except ValueError as e:
        print(f"glob pattern 錯誤：{e}")
        return
    if not results:
        print(f"找不到符合「{keyword}」的檔案")
        return
    for r in results:
        print(f"  {r}")


# ── 系統監控 ─────────────────────────────────────────

def sysinfo():
    import psutil
    cpu = psutil.cpu_percent(interval=1)
    mem = psutil.virtual_memory()
    disk = psutil.disk_usage("C:/")
    print(
        f"💻 系統狀態\n"
        f"CPU：{cpu}%\n"
        f"記憶體：{mem.percent}%（已用 {mem.used//1024//1024}MB / 共 {mem.total//1024//1024}MB）\n"
        f"磁碟 C：{disk.percent}%（已用 {disk.used//1024//1024//1024}GB / 共 {disk.total//1024//1024//1024}GB）"
    )

def process_list():
    import psutil
    print(f"{'PID':<8} {'CPU%':<7} {'記憶體MB':<10} 程序名稱")
    print("-" * 45)
    procs = sorted(psutil.process_iter(["pid","name","cpu_percent","memory_info"]),
                   key=lambda p: p.info["memory_info"].rss if p.info["memory_info"] else 0, reverse=True)
    for p in procs[:30]:
        mem_mb = p.info["memory_info"].rss // 1024 // 1024 if p.info["memory_info"] else 0
        print(f"{p.info['pid']:<8} {p.info['cpu_percent']:<7} {mem_mb:<10} {p.info['name']}")

def process_kill(name_or_pid):
    import psutil
    try:
        pid = int(name_or_pid)
        psutil.Process(pid).kill()
        print(f"✅ 已結束 PID {pid}")
    except ValueError:
        killed = 0
        for p in psutil.process_iter(["pid","name"]):
            if name_or_pid.lower() in p.info["name"].lower():
                p.kill()
                killed += 1
        print(f"✅ 已結束 {killed} 個「{name_or_pid}」程序")


# ── Windows 通知 ──────────────────────────────────────

def notify(title, message):
    try:
        from win10toast import ToastNotifier
        ToastNotifier().show_toast(title, message, duration=5, threaded=True)
        print(f"✅ 通知已送出：{title} - {message}")
    except Exception as e:
        print(f"❌ 通知失敗：{e}")


# ── 語音合成 TTS ──────────────────────────────────────

def clean_for_tts(text: str) -> str:
    """清理文字讓 TTS 更口語：去除 emoji、Markdown、特殊符號"""
    import re
    # 移除 emoji（精確範圍，避免誤刪中文字）
    emoji_pattern = re.compile(
        "["
        "\U0001F600-\U0001F64F"
        "\U0001F300-\U0001F5FF"
        "\U0001F680-\U0001F6FF"
        "\U0001F700-\U0001F9FF"
        "\U0001FA00-\U0001FAFF"
        "\u2600-\u26FF"
        "\u2700-\u27BF"
        "\u2300-\u23FF"
        "\u25A0-\u25FF"
        "\u2B00-\u2BFF"
        "\uFE00-\uFE0F"
        "\u200B-\u200D\uFEFF"
        "]+", flags=re.UNICODE
    )
    text = emoji_pattern.sub("", text)
    # 移除 Markdown
    text = re.sub(r"\*{1,3}(.*?)\*{1,3}", r"\1", text)
    text = re.sub(r"_{1,2}(.*?)_{1,2}", r"\1", text)
    text = re.sub(r"`{1,3}.*?`{1,3}", "", text, flags=re.DOTALL)
    text = re.sub(r"~~(.*?)~~", r"\1", text)
    text = re.sub(r"#{1,6}\s*", "", text)
    text = re.sub(r">\s*", "", text)
    text = re.sub(r"\[([^\]]+)\]\([^\)]+\)", r"\1", text)
    # 移除特殊符號
    text = re.sub(r"[|\\/<>{}=+\[\]~^@#$%&*_]", "", text)
    text = re.sub(r"^\s*[-•·]\s+", "", text, flags=re.MULTILINE)
    # 整理標點
    text = re.sub(r"\.{3,}", "，", text)
    lines = [l.strip() for l in text.splitlines() if l.strip()]
    text = "，".join(lines)
    text = re.sub(r"\s{2,}", " ", text).strip()
    return text

def tts(text):
    import pyttsx3
    engine = pyttsx3.init()
    engine.setProperty("rate", 180)
    engine.say(clean_for_tts(text))
    engine.runAndWait()
    print(f"✅ 已朗讀：{text}")


# ── 動作錄製/重播 ─────────────────────────────────────

_record_events = []
_record_listener = None
RECORD_DIR = Path("C:/Users/blue_/recordings")
RECORD_DIR.mkdir(exist_ok=True)

def record_start():
    from pynput import mouse, keyboard
    global _record_events, _record_listener
    _record_events = []
    start_time = time.time()

    def on_move(x, y):
        _record_events.append({"t": time.time()-start_time, "type": "move", "x": x, "y": y})
    def on_click(x, y, button, pressed):
        _record_events.append({"t": time.time()-start_time, "type": "click", "x": x, "y": y, "button": str(button), "pressed": pressed})
    def on_key(key):
        try:
            _record_events.append({"t": time.time()-start_time, "type": "key", "key": key.char})
        except AttributeError:
            _record_events.append({"t": time.time()-start_time, "type": "key", "key": str(key)})

    ml = mouse.Listener(on_move=on_move, on_click=on_click)
    kl = keyboard.Listener(on_press=on_key)
    ml.start(); kl.start()
    _record_listener = (ml, kl)
    print("✅ 開始錄製，執行 record_stop 停止")

def record_stop():
    global _record_listener
    if _record_listener:
        for l in _record_listener:
            l.stop()
        _record_listener = None
    filename = RECORD_DIR / f"rec_{dt.datetime.now().strftime('%Y%m%d_%H%M%S')}.json"
    import json
    filename.write_text(json.dumps(_record_events, ensure_ascii=False, indent=2), encoding="utf-8")
    print(f"✅ 錄製完成，已儲存：{filename}（共 {len(_record_events)} 個事件）")

def record_play(filepath):
    import json
    events = json.loads(Path(filepath).read_text(encoding="utf-8"))
    prev_t = 0
    for e in events:
        delay = e["t"] - prev_t
        if delay > 0:
            time.sleep(min(delay, 2))
        prev_t = e["t"]
        if e["type"] == "move":
            pyautogui.moveTo(e["x"], e["y"])
        elif e["type"] == "click" and e["pressed"]:
            pyautogui.click(e["x"], e["y"])
        elif e["type"] == "key":
            try:
                pyautogui.press(e["key"])
            except Exception:
                pass
    print(f"✅ 重播完成")


# ── Email ─────────────────────────────────────────────

def send_email(to: str, subject: str, body: str):
    import smtplib
    from email.mime.text import MIMEText
    from email.mime.multipart import MIMEMultipart
    smtp_host = os.getenv("SMTP_HOST", "smtp.gmail.com")
    smtp_port = int(os.getenv("SMTP_PORT", "587"))
    smtp_user = os.getenv("SMTP_USER", "")
    smtp_pass = os.getenv("SMTP_PASS", "")
    if not smtp_user or not smtp_pass:
        print("❌ 請在 .env 設定 SMTP_USER 和 SMTP_PASS")
        return
    msg = MIMEMultipart()
    msg["From"] = smtp_user
    msg["To"] = to
    msg["Subject"] = subject
    msg.attach(MIMEText(body, "plain", "utf-8"))
    with smtplib.SMTP(smtp_host, smtp_port) as s:
        s.starttls()
        s.login(smtp_user, smtp_pass)
        s.send_message(msg)
    print(f"✅ Email 已寄出到：{to}")


# ── 排程管理 ────────────────────────────────────────

SCHTASKS = "C:\\Windows\\System32\\schtasks.exe"
BOT_SCRIPT = r"C:\Users\blue_\claude-telegram-bot\bot.py"

def bot_status():
    import psutil
    bot_pids = []
    for p in psutil.process_iter(['pid', 'cmdline']):
        try:
            cmd = ' '.join(p.info.get('cmdline') or [])
            if 'bot.py' in cmd and 'psutil' not in cmd:
                bot_pids.append(p.pid)
        except Exception:
            pass
    if bot_pids:
        print(f"✅ Bot 執行中（PID: {', '.join(map(str, bot_pids))}）")
    else:
        print("❌ Bot 未執行")

def bot_restart():
    import psutil, time
    # 停掉所有 bot.py 進程
    for p in psutil.process_iter(['pid', 'cmdline']):
        try:
            cmd = ' '.join(p.info.get('cmdline') or [])
            if 'bot.py' in cmd and 'psutil' not in cmd:
                p.kill()
        except Exception:
            pass
    time.sleep(1)
    pythonw = r"C:\Users\blue_\AppData\Local\Microsoft\WindowsApps\pythonw3.12.exe"
    subprocess.Popen([pythonw, BOT_SCRIPT], cwd=str(Path(BOT_SCRIPT).parent))
    print("✅ Bot 已重啟")

def schedule_list():
    result = subprocess.run(
        [SCHTASKS, "/Query", "/FO", "CSV", "/NH"],
        capture_output=True, text=True, encoding="cp950", errors="replace"
    )
    lines = [l for l in result.stdout.strip().splitlines() if l.strip()]
    print(f"{'任務名稱':<35} {'下次執行':<25} {'狀態'}")
    print("-" * 75)
    for line in lines:
        parts = line.strip('"').split('","')
        if len(parts) >= 3:
            name = parts[0].replace("\\", "").strip()
            next_run = parts[1].strip()
            status = parts[2].strip()
            print(f"{name:<35} {next_run:<25} {status}")

def schedule_add(name: str, time_hhmm: str, script_path: str):
    # 建立排程
    subprocess.run([SCHTASKS, "/Create", "/TN", name,
                    "/TR", f"pythonw {script_path}",
                    "/SC", "DAILY", "/ST", time_hhmm, "/F"],
                   capture_output=True)
    # 設定喚醒
    ps = (
        f"$t = Get-ScheduledTask -TaskName '{name}';"
        f"$t.Settings.WakeToRun = $true;"
        f"$t.Settings.DisallowStartIfOnBatteries = $false;"
        f"$t.Settings.StopIfGoingOnBatteries = $false;"
        f"Set-ScheduledTask -TaskName '{name}' -Settings $t.Settings | Out-Null;"
        f"Write-Host '已建立'"
    )
    subprocess.run(["powershell.exe", "-Command", ps], capture_output=True, text=True)
    print(f"✅ 排程 [{name}] 已建立，每天 {time_hhmm} 執行，可喚醒電腦")

def schedule_del(name: str):
    result = subprocess.run([SCHTASKS, "/Delete", "/TN", name, "/F"],
                            capture_output=True, text=True, encoding="cp950", errors="replace")
    if result.returncode == 0:
        print(f"✅ 排程 [{name}] 已刪除")
    else:
        print(f"❌ 刪除失敗：{result.stderr.strip()}")


# ── 語音辨識 STT ─────────────────────────────────────

def stt(duration=5, file_path="", language="zh-TW"):
    """語音辨識：麥克風錄音或直接辨識音訊檔案（支援 WAV/OGG/MP3）"""
    import speech_recognition as sr, tempfile
    r = sr.Recognizer()
    try:
        if file_path:
            # 若為非 WAV 格式先用 ffmpeg 轉換
            p = Path(file_path)
            if p.suffix.lower() != ".wav":
                import imageio_ffmpeg, subprocess
                ffmpeg = imageio_ffmpeg.get_ffmpeg_exe()
                wav = tempfile.mktemp(suffix=".wav")
                subprocess.run([ffmpeg, "-y", "-i", str(p), "-ar", "16000", "-ac", "1", wav], capture_output=True)
                src_path = wav
                cleanup = True
            else:
                src_path = file_path
                cleanup = False
            with sr.AudioFile(src_path) as source:
                audio = r.record(source)
            if cleanup:
                Path(src_path).unlink(missing_ok=True)
        else:
            import sounddevice as sd, soundfile as sf
            sample_rate = 16000
            print(f"🎤 錄音 {duration} 秒，請說話...")
            recording = sd.rec(int(int(duration) * sample_rate), samplerate=sample_rate, channels=1, dtype="int16")
            sd.wait()
            with tempfile.NamedTemporaryFile(suffix=".wav", delete=False) as f:
                tmp_path = f.name
            sf.write(tmp_path, recording, sample_rate)
            with sr.AudioFile(tmp_path) as source:
                audio = r.record(source)
            Path(tmp_path).unlink(missing_ok=True)

        text = r.recognize_google(audio, language=language)
        print(f"✅ 辨識結果：{text}")
        return text
    except sr.UnknownValueError:
        print("❌ 無法辨識語音")
    except sr.RequestError as e:
        print(f"❌ STT 服務錯誤：{e}")
    except Exception as e:
        print(f"❌ STT 失敗：{e}")


# ── OCR 文字辨識 ──────────────────────────────────────

def ocr(image_path: str = ""):
    import easyocr
    reader = easyocr.Reader(["ch_tra", "en"], gpu=False)
    if image_path:
        source = image_path
    else:
        img = pyautogui.screenshot()
        source = str(SCREENSHOT_DIR / f"ocr_tmp_{dt.datetime.now().strftime('%Y%m%d_%H%M%S')}.png")
        img.save(source)
    results = reader.readtext(source)
    texts = [r[1] for r in results]
    output = "\n".join(texts)
    print(output)
    return output


# ── 自動化工作流程 ────────────────────────────────────

WORKFLOW_DIR = Path("C:/Users/blue_/workflows")
WORKFLOW_DIR.mkdir(exist_ok=True)

def workflow_save(name: str, json_content: str):
    import json
    path = WORKFLOW_DIR / f"{name}.json"
    data = json.loads(json_content)
    path.write_text(json.dumps(data, ensure_ascii=False, indent=2), encoding="utf-8")
    print(f"✅ 工作流程已儲存：{path}")

def workflow_run(json_path: str):
    import json
    path = Path(json_path)
    if not path.exists():
        path = WORKFLOW_DIR / f"{json_path}.json"
    steps = json.loads(path.read_text(encoding="utf-8"))
    print(f"▶ 執行工作流程：{path.name}（共 {len(steps)} 步）")
    for i, step in enumerate(steps, 1):
        tool = step.get("tool")
        args = step.get("args", [])
        delay = step.get("delay", 0)
        print(f"  [{i}] {tool} {args}")
        if delay:
            time.sleep(delay)
        try:
            # 動態呼叫已有工具
            tool_map = {
                "click": lambda a: pyautogui.click(int(a[0]), int(a[1])),
                "type": lambda a: pyautogui.write(" ".join(a), interval=0.05),
                "press": lambda a: pyautogui.press(a[0]),
                "hotkey": lambda a: pyautogui.hotkey(*a),
                "open": lambda a: subprocess.Popen(" ".join(a), shell=True),
                "screenshot": lambda a: screenshot(),
                "move": lambda a: pyautogui.moveTo(int(a[0]), int(a[1]), duration=0.3),
                "scroll": lambda a: pyautogui.scroll(int(a[1]) if a[0]=="up" else -int(a[1])),
                "wait": lambda a: time.sleep(float(a[0])),
                "notify": lambda a: notify(a[0], " ".join(a[1:])),
                "browser": lambda a: browser(a[0], *a[1:]),
                "file_write": lambda a: file_write(a[0], " ".join(a[1:])),
            }
            if tool in tool_map:
                tool_map[tool](args)
            else:
                print(f"    ⚠️ 未知工具：{tool}")
        except Exception as e:
            print(f"    ❌ 步驟失敗：{e}")
    print("✅ 工作流程執行完畢")


# ── 螢幕監控觸發 ──────────────────────────────────────

def screen_watch(template_path: str, command: str, interval: float = 2.0, timeout: float = 60.0):
    print(f"👁 監控中：等待 [{template_path}] 出現，超時 {timeout}s...")
    start = time.time()
    while time.time() - start < timeout:
        try:
            loc = pyautogui.locateOnScreen(template_path, confidence=0.8)
            if loc:
                print(f"✅ 偵測到目標！執行：{command}")
                subprocess.run(command, shell=True)
                return
        except Exception:
            pass
        time.sleep(interval)
    print("⏰ 監控逾時，未偵測到目標")


# ── 多螢幕支援 ───────────────────────────────────────

def monitors():
    try:
        from screeninfo import get_monitors
        for i, m in enumerate(get_monitors()):
            print(f"螢幕 {i}: {m.width}x{m.height} 位置({m.x},{m.y}) {'主螢幕' if m.is_primary else ''}")
    except Exception as e:
        print(f"❌ {e}")


# ── ZIP 壓縮 ─────────────────────────────────────────

def zip_files(source: str, dest: str):
    import zipfile
    src = Path(source)
    with zipfile.ZipFile(dest, "w", zipfile.ZIP_DEFLATED) as zf:
        if src.is_dir():
            for f in src.rglob("*"):
                zf.write(f, f.relative_to(src.parent))
        else:
            zf.write(src, src.name)
    print(f"✅ 已壓縮：{source} → {dest}")

def unzip(zip_path: str, dest: str):
    import zipfile
    with zipfile.ZipFile(zip_path, "r") as zf:
        zf.extractall(dest)
    print(f"✅ 已解壓縮：{zip_path} → {dest}")


# ── 網路下載 ─────────────────────────────────────────

def download(url: str, save_path: str = ""):
    if not save_path:
        save_path = str(SCREENSHOT_DIR / url.split("/")[-1].split("?")[0])
    r = requests.get(url, stream=True, timeout=60)
    r.raise_for_status()
    with open(save_path, "wb") as f:
        for chunk in r.iter_content(chunk_size=8192):
            f.write(chunk)
    print(f"✅ 已下載：{save_path}")


# ── 印表機 ───────────────────────────────────────────

def print_file(path: str):
    try:
        import win32api
        win32api.ShellExecute(0, "print", path, None, ".", 0)
        print(f"✅ 已送印：{path}")
    except Exception as e:
        print(f"❌ 列印失敗：{e}")


# ── WiFi 管理 ────────────────────────────────────────

def wifi_list():
    result = subprocess.run(
        ["powershell.exe", "-Command", "netsh wlan show networks mode=bssid"],
        capture_output=True, text=True, encoding="cp950", errors="replace"
    )
    print(result.stdout[:3000])

def wifi_connect(ssid: str, password: str):
    profile = f"""<?xml version="1.0"?>
<WLANProfile xmlns="http://www.microsoft.com/networking/WLAN/profile/v1">
    <name>{ssid}</name>
    <SSIDConfig><SSID><name>{ssid}</name></SSID></SSIDConfig>
    <connectionType>ESS</connectionType>
    <connectionMode>auto</connectionMode>
    <MSM><security><authEncryption>
        <authentication>WPA2PSK</authentication>
        <encryption>AES</encryption>
    </authEncryption>
    <sharedKey><keyType>passPhrase</keyType><protected>false</protected>
        <keyMaterial>{password}</keyMaterial>
    </sharedKey></security></MSM>
</WLANProfile>"""
    profile_path = Path("C:/Users/blue_/wifi_tmp.xml")
    profile_path.write_text(profile, encoding="utf-8")
    subprocess.run(["powershell.exe", "-Command",
                    f"netsh wlan add profile filename='{profile_path}'; netsh wlan connect name='{ssid}'"],
                   capture_output=True)
    profile_path.unlink(missing_ok=True)
    print(f"✅ 已嘗試連線：{ssid}")


# ── 螢幕串流 ─────────────────────────────────────────

def screen_stream(duration: int = 10, interval: float = 1.0):
    print(f"📹 開始串流 {duration} 秒（每 {interval} 秒截圖）...")
    end = time.time() + duration
    count = 0
    while time.time() < end:
        img = pyautogui.screenshot()
        filename = SCREENSHOT_DIR / f"stream_{dt.datetime.now().strftime('%H%M%S')}_{count:03d}.png"
        img.save(filename)
        count += 1
        time.sleep(interval)
    print(f"✅ 串流完成，共 {count} 張截圖存至桌面")


# ── 語音喚醒 ─────────────────────────────────────────

def wake_listen(keyword: str = "小牛馬", duration: int = 5):
    import sounddevice as sd
    import soundfile as sf
    import speech_recognition as sr
    import tempfile
    print(f"👂 監聽中，等待說出「{keyword}」...")
    sample_rate = 16000
    while True:
        recording = sd.rec(int(duration * sample_rate), samplerate=sample_rate, channels=1, dtype="int16")
        sd.wait()
        with tempfile.NamedTemporaryFile(suffix=".wav", delete=False) as f:
            tmp_path = f.name
        sf.write(tmp_path, recording, sample_rate)
        r = sr.Recognizer()
        try:
            with sr.AudioFile(tmp_path) as source:
                audio = r.record(source)
            text = r.recognize_google(audio, language="zh-TW")
            Path(tmp_path).unlink(missing_ok=True)
            if keyword in text:
                print(f"✅ 偵測到喚醒詞：{text}")
                return text
        except Exception:
            Path(tmp_path).unlink(missing_ok=True)


# ── 拖曳 / 右鍵選單 ──────────────────────────────────

def drag(x1, y1, x2, y2, duration=0.5):
    pyautogui.moveTo(int(x1), int(y1))
    pyautogui.dragTo(int(x2), int(y2), duration=float(duration), button="left")
    print(f"✅ 已拖曳 ({x1},{y1}) → ({x2},{y2})")

def right_menu(x, y, item: str = ""):
    pyautogui.rightClick(int(x), int(y))
    time.sleep(0.3)
    if item:
        import pygetwindow as gw
        pyautogui.write(item, interval=0.05)
        time.sleep(0.2)
        pyautogui.press("enter")
        print(f"✅ 已右鍵點擊並選擇：{item}")
    else:
        print(f"✅ 已右鍵點擊 ({x},{y})")


# ── 多步驟 AI 規劃 ────────────────────────────────────

def ai_plan(goal: str):
    from anthropic import Anthropic
    client = Anthropic()
    print(f"🧠 AI 規劃目標：{goal}")
    response = client.messages.create(
        model="claude-sonnet-4-6",
        max_tokens=2048,
        system="""你是一個電腦自動化規劃師。用戶給你一個目標，你要把它拆解成可執行的步驟。
每個步驟使用以下工具之一：click/type/press/hotkey/open/screenshot/browser/file_write/wait/notify
以 JSON 陣列格式回傳，例如：
[
  {"tool": "open", "args": ["notepad"], "delay": 1},
  {"tool": "type", "args": ["Hello World"]},
  {"tool": "hotkey", "args": ["ctrl","s"]}
]
只回傳 JSON，不要其他文字。""",
        messages=[{"role": "user", "content": f"目標：{goal}"}]
    )
    import json
    plan_text = response.content[0].text.strip()
    try:
        steps = json.loads(plan_text)
        print(f"📋 規劃完成，共 {len(steps)} 步：")
        for i, s in enumerate(steps, 1):
            print(f"  [{i}] {s.get('tool')} {s.get('args', [])}")
        print("\n▶ 開始執行...")
        tool_map = {
            "click": lambda a: pyautogui.click(int(a[0]), int(a[1])),
            "type": lambda a: pyautogui.write(" ".join(str(x) for x in a), interval=0.05),
            "press": lambda a: pyautogui.press(a[0]),
            "hotkey": lambda a: pyautogui.hotkey(*a),
            "open": lambda a: subprocess.Popen(" ".join(str(x) for x in a), shell=True),
            "screenshot": lambda a: screenshot(),
            "wait": lambda a: time.sleep(float(a[0])),
            "notify": lambda a: notify(a[0], a[1] if len(a) > 1 else ""),
            "move": lambda a: pyautogui.moveTo(int(a[0]), int(a[1]), duration=0.3),
            "scroll": lambda a: pyautogui.scroll(int(a[1]) if a[0] == "up" else -int(a[1])),
            "file_write": lambda a: file_write(a[0], " ".join(str(x) for x in a[1:])),
            "browser": lambda a: browser(a[0], *a[1:]),
        }
        for i, step in enumerate(steps, 1):
            t = step.get("tool")
            a = step.get("args", [])
            d = step.get("delay", 0)
            if d: time.sleep(d)
            if t in tool_map:
                tool_map[t](a)
                print(f"  ✅ 步驟 {i} 完成")
            else:
                print(f"  ⚠️ 未知工具：{t}")
        print("✅ 全部執行完畢")
    except json.JSONDecodeError:
        print(f"規劃結果：\n{plan_text}")


# ── 剪貼簿歷史 ───────────────────────────────────────

_clipboard_history = []

def clipboard_history():
    import pyperclip
    current = pyperclip.paste()
    if current and (not _clipboard_history or _clipboard_history[-1] != current):
        _clipboard_history.append(current)
    if not _clipboard_history:
        print("剪貼簿歷史是空的")
        return
    for i, item in enumerate(_clipboard_history[-20:], 1):
        preview = item[:80].replace("\n", "↵")
        print(f"  [{i}] {preview}")


# ── 虛擬桌面 ─────────────────────────────────────────

def vdesktop(action: str):
    if action == "left":
        pyautogui.hotkey("ctrl", "win", "left")
        print("✅ 切換到左邊虛擬桌面")
    elif action == "right":
        pyautogui.hotkey("ctrl", "win", "right")
        print("✅ 切換到右邊虛擬桌面")
    elif action == "new":
        pyautogui.hotkey("ctrl", "win", "d")
        print("✅ 已建立新虛擬桌面")
    else:
        print(f"❌ 未知動作：{action}（可用：left/right/new）")


# ── 電源管理 ─────────────────────────────────────────

def power(action: str):
    cmds = {
        "sleep":    "rundll32.exe powrprof.dll,SetSuspendState 0,1,0",
        "restart":  "shutdown /r /t 5",
        "shutdown": "shutdown /s /t 5",
    }
    if action not in cmds:
        print(f"❌ 未知動作：{action}（可用：sleep/restart/shutdown）")
        return
    subprocess.run(["powershell.exe", "-Command", cmds[action]])
    print(f"✅ 已執行：{action}")


# ── 藍牙管理 ─────────────────────────────────────────

def bt_scan():
    try:
        import asyncio
        import bleak
        async def _scan():
            devices = await bleak.BleakScanner.discover(timeout=5.0)
            return devices
        devices = asyncio.run(_scan())
        if not devices:
            print("找不到藍牙裝置")
            return
        for d in devices:
            print(f"  {d.address}  {d.name or '(未知)'}")
    except Exception as e:
        print(f"❌ 藍牙掃描失敗：{e}")

def bt_connect(mac: str):
    try:
        import asyncio
        import bleak
        async def _connect():
            async with bleak.BleakClient(mac) as client:
                print(f"✅ 已連線：{mac}（服務數：{len(client.services)}）")
                await asyncio.sleep(3)
        asyncio.run(_connect())
    except Exception as e:
        print(f"❌ 連線失敗：{e}")


# ── 程式碼執行 ───────────────────────────────────────

def run_python(code: str):
    import traceback, contextlib
    buf = io.StringIO()
    try:
        with contextlib.redirect_stdout(buf), contextlib.redirect_stderr(buf):
            exec(code, {"__builtins__": __builtins__})
        result = buf.getvalue()
        print(result if result else "✅ 執行完成（無輸出）")
    except Exception:
        print(f"❌ 執行錯誤：\n{traceback.format_exc()}")

def run_shell(cmd: str):
    result = subprocess.run(
        ["powershell.exe", "-Command", cmd],
        capture_output=True, text=True, encoding="cp950", errors="replace"
    )
    output = result.stdout + result.stderr
    print(output.strip() or "✅ 執行完成")


# ── 文件處理 ─────────────────────────────────────────

def word_read(path: str):
    from docx import Document
    doc = Document(path)
    text = "\n".join(p.text for p in doc.paragraphs if p.text.strip())
    print(text[:3000])

def word_write(path: str, content: str):
    from docx import Document
    doc = Document()
    for line in content.split("\n"):
        doc.add_paragraph(line)
    doc.save(path)
    print(f"✅ 已寫入：{path}")

def excel_read(path: str, sheet: str = ""):
    import openpyxl
    wb = openpyxl.load_workbook(path, data_only=True)
    ws = wb[sheet] if sheet and sheet in wb.sheetnames else wb.active
    for row in ws.iter_rows(values_only=True):
        print("\t".join(str(c) if c is not None else "" for c in row))

def excel_write(path: str, sheet: str, data_json: str):
    import openpyxl, json
    data = json.loads(data_json)
    try:
        wb = openpyxl.load_workbook(path)
    except FileNotFoundError:
        wb = openpyxl.Workbook()
    ws = wb[sheet] if sheet in wb.sheetnames else wb.create_sheet(sheet)
    for row in data:
        ws.append(row)
    wb.save(path)
    print(f"✅ 已寫入 Excel：{path} [{sheet}]")

def pdf_read(path: str):
    import fitz
    doc = fitz.open(path)
    text = "\n".join(page.get_text() for page in doc)
    print(text[:3000])


# ── 螢幕變化偵測 ──────────────────────────────────────

def screen_diff(interval: float = 1.0, duration: float = 30.0, region=None):
    import numpy as np
    import cv2
    print(f"👁 監控螢幕變化（{duration}秒）...")
    prev = None
    end = time.time() + duration
    changes = 0
    while time.time() < end:
        img = pyautogui.screenshot(region=region)
        arr = np.array(img.convert("L"))
        if prev is not None:
            diff = cv2.absdiff(prev, arr)
            score = diff.mean()
            if score > 5:
                changes += 1
                ts = dt.datetime.now().strftime("%H:%M:%S")
                print(f"  [{ts}] 偵測到變化！差異分數：{score:.1f}")
        prev = arr
        time.sleep(interval)
    print(f"✅ 監控結束，共偵測到 {changes} 次變化")


# ── 網頁爬蟲 ─────────────────────────────────────────

def scrape(url: str, selector: str = "body"):
    from bs4 import BeautifulSoup
    res = requests.get(url, headers={"User-Agent": "Mozilla/5.0"}, timeout=15)
    res.raise_for_status()
    soup = BeautifulSoup(res.text, "html.parser")
    elements = soup.select(selector)
    if not elements:
        print(f"找不到選擇器：{selector}")
        return
    for el in elements[:10]:
        text = el.get_text(strip=True)
        if text:
            print(text[:200])


# ── 圖片編輯 ─────────────────────────────────────────

def img_edit(action: str, path: str, *args):
    from PIL import Image, ImageDraw, ImageFont
    img = Image.open(path)

    if action == "crop":
        x1, y1, x2, y2 = int(args[0]), int(args[1]), int(args[2]), int(args[3])
        img = img.crop((x1, y1, x2, y2))
    elif action == "resize":
        w, h = int(args[0]), int(args[1])
        img = img.resize((w, h), Image.LANCZOS)
    elif action == "text":
        text = args[0]
        x, y = int(args[1]) if len(args) > 1 else img.width // 2, int(args[2]) if len(args) > 2 else img.height - 50
        draw = ImageDraw.Draw(img)
        try:
            font = ImageFont.truetype("C:/Windows/Fonts/msjhbd.ttc", 36)
        except Exception:
            font = ImageFont.load_default()
        draw.text((x, y), text, font=font, fill=(255, 255, 255))
    elif action == "merge":
        img2 = Image.open(args[0])
        new = Image.new("RGB", (img.width + img2.width, max(img.height, img2.height)))
        new.paste(img, (0, 0))
        new.paste(img2, (img.width, 0))
        img = new
    elif action == "grayscale":
        img = img.convert("L")

    out_path = args[-1] if args and str(args[-1]).endswith((".png",".jpg",".jpeg")) else path
    img.save(out_path)
    print(f"✅ 圖片已儲存：{out_path}")


# ── Google Drive ──────────────────────────────────────

def gdrive_upload(file_path: str, folder_id: str = ""):
    try:
        from googleapiclient.discovery import build
        from googleapiclient.http import MediaFileUpload
        from google.oauth2.credentials import Credentials
        creds_path = Path("C:/Users/blue_/gdrive_token.json")
        if not creds_path.exists():
            print("❌ 找不到 C:/Users/blue_/gdrive_token.json，請先完成 Google 授權")
            return
        creds = Credentials.from_authorized_user_file(str(creds_path))
        service = build("drive", "v3", credentials=creds)
        meta = {"name": Path(file_path).name}
        if folder_id:
            meta["parents"] = [folder_id]
        media = MediaFileUpload(file_path)
        f = service.files().create(body=meta, media_body=media, fields="id").execute()
        print(f"✅ 已上傳，檔案 ID：{f.get('id')}")
    except Exception as e:
        print(f"❌ 上傳失敗：{e}")

def gdrive_download(file_id: str, save_path: str):
    try:
        from googleapiclient.discovery import build
        from googleapiclient.http import MediaIoBaseDownload
        from google.oauth2.credentials import Credentials
        creds = Credentials.from_authorized_user_file("C:/Users/blue_/gdrive_token.json")
        service = build("drive", "v3", credentials=creds)
        req = service.files().get_media(fileId=file_id)
        with open(save_path, "wb") as f:
            downloader = MediaIoBaseDownload(f, req)
            done = False
            while not done:
                _, done = downloader.next_chunk()
        print(f"✅ 已下載：{save_path}")
    except Exception as e:
        print(f"❌ 下載失敗：{e}")


# ── 資料庫操作 ───────────────────────────────────────

def db_query(db_path: str, sql: str):
    conn = sqlite3.connect(db_path)
    try:
        cur = conn.execute(sql)
        if cur.description:
            headers = [d[0] for d in cur.description]
            print("\t".join(headers))
            print("-" * 60)
            for row in cur.fetchall():
                print("\t".join(str(c) for c in row))
        else:
            conn.commit()
            print(f"✅ 執行成功，影響 {cur.rowcount} 列")
    finally:
        conn.close()

def db_mysql(host: str, database: str, sql: str, user: str = "root", password: str = ""):
    import pymysql
    conn = pymysql.connect(host=host, user=user, password=password, database=database, charset="utf8mb4")
    try:
        with conn.cursor() as cur:
            cur.execute(sql)
            if cur.description:
                headers = [d[0] for d in cur.description]
                print("\t".join(headers))
                for row in cur.fetchall():
                    print("\t".join(str(c) for c in row))
            else:
                conn.commit()
                print(f"✅ 執行成功，影響 {cur.rowcount} 列")
    finally:
        conn.close()


# ── 加密/解密 ────────────────────────────────────────

def encrypt(file_path: str, password: str):
    from cryptography.fernet import Fernet
    from cryptography.hazmat.primitives.kdf.pbkdf2 import PBKDF2HMAC
    from cryptography.hazmat.primitives import hashes
    import base64
    salt = b"xinyang501_salt_"
    kdf = PBKDF2HMAC(algorithm=hashes.SHA256(), length=32, salt=salt, iterations=100000)
    key = base64.urlsafe_b64encode(kdf.derive(password.encode()))
    f = Fernet(key)
    data = Path(file_path).read_bytes()
    out_path = file_path + ".enc"
    Path(out_path).write_bytes(f.encrypt(data))
    print(f"✅ 已加密：{out_path}")

def decrypt(file_path: str, password: str):
    from cryptography.fernet import Fernet
    from cryptography.hazmat.primitives.kdf.pbkdf2 import PBKDF2HMAC
    from cryptography.hazmat.primitives import hashes
    import base64
    salt = b"xinyang501_salt_"
    kdf = PBKDF2HMAC(algorithm=hashes.SHA256(), length=32, salt=salt, iterations=100000)
    key = base64.urlsafe_b64encode(kdf.derive(password.encode()))
    f = Fernet(key)
    data = Path(file_path).read_bytes()
    out_path = file_path.replace(".enc", ".dec")
    Path(out_path).write_bytes(f.decrypt(data))
    print(f"✅ 已解密：{out_path}")


# ── 剪貼簿監控 ───────────────────────────────────────

def clipboard_watch(duration: float = 30.0):
    import pyperclip
    print(f"📋 監控剪貼簿 {duration} 秒...")
    prev = pyperclip.paste()
    end = time.time() + duration
    changes = []
    while time.time() < end:
        current = pyperclip.paste()
        if current != prev and current:
            ts = dt.datetime.now().strftime("%H:%M:%S")
            preview = current[:80].replace("\n", "↵")
            print(f"  [{ts}] 新內容：{preview}")
            changes.append(current)
            prev = current
        time.sleep(0.5)
    print(f"✅ 監控結束，共偵測到 {len(changes)} 次變化")


# ── QR Code ──────────────────────────────────────────

def qr_gen(content: str, save_path: str = ""):
    import qrcode
    qr = qrcode.make(content)
    if not save_path:
        save_path = str(SCREENSHOT_DIR / f"qr_{dt.datetime.now().strftime('%Y%m%d_%H%M%S')}.png")
    qr.save(save_path)
    print(f"✅ QR Code 已生成：{save_path}")

def qr_scan(image_path: str = ""):
    try:
        from pyzbar.pyzbar import decode
        from PIL import Image
        if not image_path:
            img = pyautogui.screenshot()
        else:
            img = Image.open(image_path)
        results = decode(img)
        if not results:
            print("❌ 未偵測到 QR Code")
            return
        for r in results:
            print(f"✅ 掃描結果：{r.data.decode('utf-8')}")
    except Exception as e:
        print(f"❌ 掃描失敗：{e}")


# ── 螢幕錄影 ────────────────────────────────────────

def screen_record(duration: float = 10.0, output: str = ""):
    try:
        import mss, cv2, numpy as np, time as t
        out_path = output or str(Path.home() / "Desktop" / f"record_{dt.datetime.now().strftime('%H%M%S')}.mp4")
        with mss.mss() as sct:
            mon = sct.monitors[1]
            w, h = mon["width"], mon["height"]
            fps = 10
            writer = cv2.VideoWriter(out_path, cv2.VideoWriter_fourcc(*"mp4v"), fps, (w, h))
            end = t.time() + duration
            while t.time() < end:
                frame = np.array(sct.grab(mon))
                frame = cv2.cvtColor(frame, cv2.COLOR_BGRA2BGR)
                writer.write(frame)
                t.sleep(1 / fps)
            writer.release()
        print(f"✅ 錄影完成：{out_path}")
    except Exception as e:
        print(f"❌ 錄影失敗：{e}")


def webcam_capture(output: str = ""):
    try:
        import cv2
        out_path = output or str(Path.home() / "Desktop" / f"webcam_{dt.datetime.now().strftime('%H%M%S')}.jpg")
        cap = cv2.VideoCapture(0)
        ret, frame = cap.read()
        cap.release()
        if ret:
            cv2.imwrite(out_path, frame)
            print(f"✅ 已拍照：{out_path}")
        else:
            print("❌ 無法存取攝影機")
    except Exception as e:
        print(f"❌ 攝影機失敗：{e}")


# ── 翻譯 ────────────────────────────────────────────

def translate(text: str, target: str = "zh-TW", source: str = "auto"):
    try:
        from deep_translator import GoogleTranslator
        result = GoogleTranslator(source=source, target=target).translate(text)
        print(result)
    except Exception as e:
        print(f"❌ 翻譯失敗：{e}")


# ── 資料視覺化 ───────────────────────────────────────

def chart(chart_type: str, data_json: str, title: str = "", output: str = ""):
    try:
        import matplotlib
        matplotlib.use("Agg")
        import matplotlib.pyplot as plt
        import json
        data = json.loads(data_json)
        out_path = output or str(Path.home() / "Desktop" / f"chart_{dt.datetime.now().strftime('%H%M%S')}.png")
        fig, ax = plt.subplots()
        if chart_type == "line":
            for label, values in data.items():
                ax.plot(values, label=label)
            ax.legend()
        elif chart_type == "bar":
            ax.bar(list(data.keys()), list(data.values()))
        elif chart_type == "pie":
            ax.pie(list(data.values()), labels=list(data.keys()), autopct="%1.1f%%")
        if title:
            ax.set_title(title)
        plt.tight_layout()
        plt.savefig(out_path)
        plt.close()
        print(f"✅ 圖表已存：{out_path}")
    except Exception as e:
        print(f"❌ 圖表生成失敗：{e}")


# ── PowerPoint ───────────────────────────────────────

def pptx_read(path: str):
    try:
        from pptx import Presentation
        prs = Presentation(path)
        for i, slide in enumerate(prs.slides, 1):
            texts = [sh.text for sh in slide.shapes if sh.has_text_frame]
            print(f"[投影片 {i}] " + " | ".join(t for t in texts if t.strip()))
    except Exception as e:
        print(f"❌ 讀取失敗：{e}")

def pptx_create(path: str, slides_json: str):
    try:
        from pptx import Presentation
        from pptx.util import Inches, Pt
        import json
        slides = json.loads(slides_json)
        prs = Presentation()
        for s in slides:
            layout = prs.slide_layouts[1]
            slide = prs.slides.add_slide(layout)
            if "title" in s:
                slide.shapes.title.text = s["title"]
            if "body" in s and slide.placeholders[1]:
                slide.placeholders[1].text = s["body"]
        prs.save(path)
        print(f"✅ 已建立簡報：{path}")
    except Exception as e:
        print(f"❌ 建立失敗：{e}")


# ── REST API ─────────────────────────────────────────

def api_call(method: str, url: str, headers_json: str = "{}", body_json: str = "{}"):
    try:
        import json
        headers = json.loads(headers_json)
        body = json.loads(body_json)
        resp = requests.request(method.upper(), url, headers=headers, json=body if body else None, timeout=30)
        try:
            print(json.dumps(resp.json(), ensure_ascii=False, indent=2)[:3000])
        except Exception:
            print(resp.text[:3000])
    except Exception as e:
        print(f"❌ API 呼叫失敗：{e}")


# ── 程序守護 ─────────────────────────────────────────

def watchdog(process_name: str, script: str, duration: float = 60.0):
    try:
        import psutil, time as t
        print(f"🐕 開始監控 [{process_name}]，持續 {duration} 秒...")
        end = t.time() + duration
        restarts = 0
        while t.time() < end:
            running = any(p.name().lower() == process_name.lower() for p in psutil.process_iter())
            if not running:
                subprocess.Popen(["pythonw" if script.endswith(".py") else script, script] if script.endswith(".py") else [script])
                restarts += 1
                print(f"⚠️ [{process_name}] 已重啟（第 {restarts} 次）")
            t.sleep(5)
        print(f"✅ 守護結束，共重啟 {restarts} 次")
    except Exception as e:
        print(f"❌ 守護失敗：{e}")


# ── SSH ──────────────────────────────────────────────

def ssh_run(host: str, user: str, password: str, command: str, port: int = 22):
    try:
        import paramiko
        client = paramiko.SSHClient()
        client.set_missing_host_key_policy(paramiko.AutoAddPolicy())
        client.connect(host, port=port, username=user, password=password, timeout=15)
        _, stdout, stderr = client.exec_command(command)
        out = stdout.read().decode(errors="replace")
        err = stderr.read().decode(errors="replace")
        client.close()
        print((out + err).strip() or "（執行完畢，無輸出）")
    except Exception as e:
        print(f"❌ SSH 失敗：{e}")

def sftp_upload(host: str, user: str, password: str, local: str, remote: str, port: int = 22):
    try:
        import paramiko
        t = paramiko.Transport((host, port))
        t.connect(username=user, password=password)
        sftp = paramiko.SFTPClient.from_transport(t)
        sftp.put(local, remote)
        sftp.close(); t.close()
        print(f"✅ 已上傳：{local} → {remote}")
    except Exception as e:
        print(f"❌ SFTP 上傳失敗：{e}")

def sftp_download(host: str, user: str, password: str, remote: str, local: str, port: int = 22):
    try:
        import paramiko
        t = paramiko.Transport((host, port))
        t.connect(username=user, password=password)
        sftp = paramiko.SFTPClient.from_transport(t)
        sftp.get(remote, local)
        sftp.close(); t.close()
        print(f"✅ 已下載：{remote} → {local}")
    except Exception as e:
        print(f"❌ SFTP 下載失敗：{e}")


# ── 網路診斷 ─────────────────────────────────────────

def net_ping(host: str, count: int = 4):
    result = subprocess.run(["ping", "-n", str(count), host], capture_output=True, text=True, encoding="cp950", errors="replace")
    print(result.stdout.strip())

def net_traceroute(host: str):
    result = subprocess.run(["tracert", host], capture_output=True, text=True, encoding="cp950", errors="replace", timeout=60)
    print(result.stdout[:3000])

def net_portscan(host: str, ports: str = "22,80,443,3306,3389,8080"):
    import socket
    results = []
    for p in [int(x) for x in ports.split(",")]:
        try:
            s = socket.socket()
            s.settimeout(1)
            r = s.connect_ex((host, p))
            results.append(f"Port {p}: {'開放 ✅' if r == 0 else '關閉 ❌'}")
            s.close()
        except Exception:
            results.append(f"Port {p}: 錯誤")
    print("\n".join(results))


# ── Windows 服務 ─────────────────────────────────────

def win_service(action: str, name: str = ""):
    try:
        if action == "list":
            r = subprocess.run(["powershell.exe", "-Command",
                "Get-Service | Select-Object Name,Status | Format-Table -AutoSize"],
                capture_output=True, text=True, encoding="utf-8", errors="replace")
            print(r.stdout[:3000])
        elif action in ("start", "stop"):
            cmd = f"{'Start' if action=='start' else 'Stop'}-Service -Name '{name}' -Force"
            r = subprocess.run(["powershell.exe", "-Command", cmd],
                capture_output=True, text=True, encoding="utf-8", errors="replace")
            print(r.stdout or r.stderr or f"✅ {action} {name}")
    except Exception as e:
        print(f"❌ 服務操作失敗：{e}")


# ── PDF 編輯 ─────────────────────────────────────────

def pdf_merge(paths_json: str, output: str):
    try:
        import fitz, json
        paths = json.loads(paths_json)
        writer = fitz.open()
        for p in paths:
            writer.insert_pdf(fitz.open(p))
        writer.save(output)
        print(f"✅ 已合併 {len(paths)} 個 PDF：{output}")
    except Exception as e:
        print(f"❌ 合併失敗：{e}")

def pdf_split(path: str, output_dir: str):
    try:
        import fitz
        doc = fitz.open(path)
        Path(output_dir).mkdir(parents=True, exist_ok=True)
        for i, page in enumerate(doc):
            out = fitz.open()
            out.insert_pdf(doc, from_page=i, to_page=i)
            out.save(str(Path(output_dir) / f"page_{i+1}.pdf"))
        print(f"✅ 已分割 {len(doc)} 頁到：{output_dir}")
    except Exception as e:
        print(f"❌ 分割失敗：{e}")

def pdf_watermark(path: str, text: str, output: str = ""):
    try:
        import fitz
        doc = fitz.open(path)
        out_path = output or path.replace(".pdf", "_wm.pdf")
        for page in doc:
            page.insert_text((page.rect.width/2 - 50, page.rect.height/2),
                text, fontsize=40, color=(0.8, 0.8, 0.8), rotate=45)
        doc.save(out_path)
        print(f"✅ 已加浮水印：{out_path}")
    except Exception as e:
        print(f"❌ 浮水印失敗：{e}")


# ── 音訊處理 ─────────────────────────────────────────

def audio_convert(input_path: str, output_path: str):
    try:
        from pydub import AudioSegment
        fmt = Path(output_path).suffix.lstrip(".")
        AudioSegment.from_file(input_path).export(output_path, format=fmt)
        print(f"✅ 已轉換：{output_path}")
    except Exception as e:
        print(f"❌ 音訊轉換失敗：{e}")

def audio_trim(input_path: str, start_ms: int, end_ms: int, output_path: str = ""):
    try:
        from pydub import AudioSegment
        audio = AudioSegment.from_file(input_path)[start_ms:end_ms]
        out = output_path or input_path.replace(".", "_trim.")
        audio.export(out, format=Path(out).suffix.lstrip("."))
        print(f"✅ 已剪輯：{out}")
    except Exception as e:
        print(f"❌ 音訊剪輯失敗：{e}")


# ── 推播通知 ─────────────────────────────────────────

def discord_notify(webhook_url: str, message: str):
    try:
        resp = requests.post(webhook_url, json={"content": message}, timeout=10)
        print(f"✅ Discord 已發送（{resp.status_code}）")
    except Exception as e:
        print(f"❌ Discord 發送失敗：{e}")

def line_notify(token: str, message: str):
    try:
        resp = requests.post(
            "https://notify-api.line.me/api/notify",
            headers={"Authorization": f"Bearer {token}"},
            data={"message": message}, timeout=10
        )
        print(f"✅ LINE 已發送（{resp.status_code}）")
    except Exception as e:
        print(f"❌ LINE 發送失敗：{e}")


# ── 磁碟清理 ─────────────────────────────────────────

def disk_clean(action: str = "list"):
    try:
        import tempfile, shutil
        tmp = Path(tempfile.gettempdir())
        if action == "list":
            files = list(tmp.rglob("*"))
            total = sum(f.stat().st_size for f in files if f.is_file())
            print(f"暫存資料夾：{tmp}\n檔案數：{len(files)}\n佔用空間：{total/1024/1024:.1f} MB")
        elif action == "clean":
            count = 0
            for f in tmp.iterdir():
                try:
                    if f.is_file():
                        f.unlink(); count += 1
                    elif f.is_dir():
                        shutil.rmtree(f, ignore_errors=True); count += 1
                except Exception:
                    pass
            print(f"✅ 已清理 {count} 個暫存項目")
    except Exception as e:
        print(f"❌ 磁碟清理失敗：{e}")

def backup(src: str, dest: str):
    try:
        import shutil
        out = Path(dest) / f"{Path(src).name}_{dt.datetime.now().strftime('%Y%m%d_%H%M%S')}.zip"
        shutil.make_archive(str(out).replace(".zip",""), "zip", src)
        print(f"✅ 備份完成：{out}")
    except Exception as e:
        print(f"❌ 備份失敗：{e}")


# ── Windows 登錄檔 ────────────────────────────────────

def registry_read(key_path: str, value_name: str = ""):
    try:
        import winreg
        parts = key_path.split("\\", 1)
        roots = {"HKEY_LOCAL_MACHINE": winreg.HKEY_LOCAL_MACHINE,
                 "HKEY_CURRENT_USER": winreg.HKEY_CURRENT_USER,
                 "HKLM": winreg.HKEY_LOCAL_MACHINE, "HKCU": winreg.HKEY_CURRENT_USER}
        root = roots[parts[0]]
        with winreg.OpenKey(root, parts[1]) as k:
            if value_name:
                val, _ = winreg.QueryValueEx(k, value_name)
                print(f"{value_name} = {val}")
            else:
                i = 0
                while True:
                    try:
                        n, v, _ = winreg.EnumValue(k, i)
                        print(f"{n} = {v}")
                        i += 1
                    except OSError:
                        break
    except Exception as e:
        print(f"❌ 讀取登錄檔失敗：{e}")

def registry_write(key_path: str, value_name: str, value: str):
    try:
        import winreg
        parts = key_path.split("\\", 1)
        roots = {"HKEY_LOCAL_MACHINE": winreg.HKEY_LOCAL_MACHINE,
                 "HKEY_CURRENT_USER": winreg.HKEY_CURRENT_USER,
                 "HKLM": winreg.HKEY_LOCAL_MACHINE, "HKCU": winreg.HKEY_CURRENT_USER}
        root = roots[parts[0]]
        with winreg.OpenKey(root, parts[1], 0, winreg.KEY_SET_VALUE) as k:
            winreg.SetValueEx(k, value_name, 0, winreg.REG_SZ, value)
        print(f"✅ 已寫入：{value_name} = {value}")
    except Exception as e:
        print(f"❌ 寫入登錄檔失敗：{e}")


# ── 影片處理 ─────────────────────────────────────────

def video_screenshot(path: str, second: float = 0, output: str = ""):
    try:
        import cv2
        cap = cv2.VideoCapture(path)
        cap.set(cv2.CAP_PROP_POS_MSEC, second * 1000)
        ret, frame = cap.read()
        cap.release()
        out = output or path.replace(".mp4", f"_frame{int(second)}s.jpg")
        if ret:
            cv2.imwrite(out, frame)
            print(f"✅ 已擷取畫面：{out}")
        else:
            print("❌ 無法讀取影片")
    except Exception as e:
        print(f"❌ 影片截圖失敗：{e}")

def video_trim(path: str, start_sec: float, end_sec: float, output: str = ""):
    try:
        out = output or path.replace(".mp4", f"_trim.mp4")
        subprocess.run([
            "ffmpeg", "-y", "-i", path,
            "-ss", str(start_sec), "-to", str(end_sec),
            "-c", "copy", out
        ], capture_output=True)
        print(f"✅ 已剪輯：{out}")
    except Exception as e:
        print(f"❌ 影片剪輯失敗：{e}")


# ── 多螢幕管理 ───────────────────────────────────────

def monitor_list():
    try:
        from screeninfo import get_monitors
        for m in get_monitors():
            print(f"{'主螢幕 ' if m.is_primary else '副螢幕 '}{m.width}x{m.height} @({m.x},{m.y}) name={m.name}")
    except Exception as e:
        print(f"❌ 取得螢幕資訊失敗：{e}")


# ── Email 讀取 (IMAP) ────────────────────────────────

def email_read(host: str, user: str, password: str, folder: str = "INBOX", count: int = 5):
    try:
        import imapclient, email as _email
        from email.header import decode_header
        client = imapclient.IMAPClient(host, ssl=True)
        client.login(user, password)
        client.select_folder(folder)
        msgs = client.search(["ALL"])
        recent = msgs[-count:] if len(msgs) >= count else msgs
        results = []
        for uid in reversed(recent):
            raw = client.fetch([uid], ["RFC822"])[uid][b"RFC822"]
            msg = _email.message_from_bytes(raw)
            subj_raw, enc = decode_header(msg["Subject"])[0]
            subject = subj_raw.decode(enc or "utf-8") if isinstance(subj_raw, bytes) else subj_raw
            sender = msg["From"]
            date = msg["Date"]
            body = ""
            if msg.is_multipart():
                for part in msg.walk():
                    if part.get_content_type() == "text/plain":
                        body = part.get_payload(decode=True).decode(errors="replace")[:200]
                        break
            else:
                body = msg.get_payload(decode=True).decode(errors="replace")[:200]
            results.append(f"[{date}]\n寄件人：{sender}\n主旨：{subject}\n{body}\n{'─'*30}")
        client.logout()
        print("\n".join(results))
    except Exception as e:
        print(f"❌ 讀取郵件失敗：{e}")


# ── Google Calendar ───────────────────────────────────

def gcal_list(days: int = 7):
    try:
        from google.oauth2.credentials import Credentials
        from googleapiclient.discovery import build
        from datetime import timezone, timedelta
        creds_path = Path("C:/Users/blue_/claude-telegram-bot/gcal_token.json")
        if not creds_path.exists():
            print("❌ 未找到 Google Calendar 憑證（gcal_token.json）")
            return
        creds = Credentials.from_authorized_user_file(str(creds_path))
        service = build("calendar", "v3", credentials=creds)
        now = dt.datetime.now(dt.timezone.utc)
        end = now + timedelta(days=days)
        events = service.events().list(
            calendarId="primary",
            timeMin=now.isoformat(),
            timeMax=end.isoformat(),
            maxResults=20, singleEvents=True, orderBy="startTime"
        ).execute().get("items", [])
        if not events:
            print(f"未來 {days} 天沒有行程")
            return
        for e in events:
            start = e["start"].get("dateTime", e["start"].get("date"))
            print(f"📅 {start}  {e.get('summary','（無標題）')}")
    except Exception as e:
        print(f"❌ 讀取行事曆失敗：{e}")

def gcal_add(title: str, start: str, end: str, description: str = ""):
    try:
        from google.oauth2.credentials import Credentials
        from googleapiclient.discovery import build
        creds_path = Path("C:/Users/blue_/claude-telegram-bot/gcal_token.json")
        if not creds_path.exists():
            print("❌ 未找到 Google Calendar 憑證（gcal_token.json）")
            return
        creds = Credentials.from_authorized_user_file(str(creds_path))
        service = build("calendar", "v3", credentials=creds)
        event = {
            "summary": title,
            "description": description,
            "start": {"dateTime": start, "timeZone": "Asia/Taipei"},
            "end": {"dateTime": end, "timeZone": "Asia/Taipei"},
        }
        created = service.events().insert(calendarId="primary", body=event).execute()
        print(f"✅ 行程已新增：{created.get('htmlLink')}")
    except Exception as e:
        print(f"❌ 新增行程失敗：{e}")


# ── 全域快捷鍵 ───────────────────────────────────────

def global_hotkey_listen(hotkey: str, command: str, duration: float = 60.0):
    try:
        import keyboard as kb, time as t
        triggered = []
        def on_trigger():
            triggered.append(dt.datetime.now().strftime("%H:%M:%S"))
            subprocess.run(command, shell=True)
        kb.add_hotkey(hotkey, on_trigger)
        print(f"🎹 監聽快捷鍵 [{hotkey}]，持續 {duration} 秒...")
        t.sleep(duration)
        kb.remove_all_hotkeys()
        print(f"✅ 共觸發 {len(triggered)} 次：{triggered}")
    except Exception as e:
        print(f"❌ 快捷鍵監聽失敗：{e}")


# ── Git 操作 ─────────────────────────────────────────

def git_op(action: str, repo: str = ".", message: str = "", remote: str = "origin", branch: str = "master"):
    try:
        import git as _git
        repo_obj = _git.Repo(repo)
        if action == "status":
            print(repo_obj.git.status())
        elif action == "log":
            for c in list(repo_obj.iter_commits())[:10]:
                print(f"{c.hexsha[:7]} [{c.authored_datetime.strftime('%Y-%m-%d %H:%M')}] {c.message.strip()[:60]}")
        elif action == "pull":
            result = repo_obj.remotes[remote].pull()
            print(f"✅ Pull 完成：{result[0].commit.hexsha[:7]}")
        elif action == "add":
            repo_obj.git.add(A=True)
            print("✅ 已 git add -A")
        elif action == "commit":
            repo_obj.index.commit(message or "auto commit")
            print(f"✅ 已 commit：{message}")
        elif action == "push":
            repo_obj.remotes[remote].push(branch)
            print(f"✅ 已 push 到 {remote}/{branch}")
        elif action == "diff":
            print(repo_obj.git.diff()[:3000] or "（無變更）")
    except Exception as e:
        print(f"❌ Git 操作失敗：{e}")


# ── 硬體監控進階 ─────────────────────────────────────

def hw_monitor():
    try:
        import psutil
        cpu_temp = ""
        try:
            temps = psutil.sensors_temperatures()
            if temps:
                for name, entries in temps.items():
                    for e in entries[:2]:
                        cpu_temp += f"{name}: {e.current}°C  "
        except Exception:
            cpu_temp = "（不支援溫度感測）"

        battery = psutil.sensors_battery()
        bat_str = f"{battery.percent:.0f}% {'充電中' if battery.power_plugged else '使用電池'}" if battery else "無電池"

        gpu_str = ""
        try:
            import GPUtil
            gpus = GPUtil.getGPUs()
            for g in gpus:
                gpu_str += f"\nGPU [{g.name}] 使用率：{g.load*100:.0f}% | 記憶體：{g.memoryUsed:.0f}/{g.memoryTotal:.0f}MB | 溫度：{g.temperature}°C"
        except Exception:
            gpu_str = "\nGPU：（未偵測到 NVIDIA GPU）"

        print(
            f"🌡 溫度：{cpu_temp}\n"
            f"🔋 電池：{bat_str}"
            f"{gpu_str}"
        )
    except Exception as e:
        print(f"❌ 硬體監控失敗：{e}")


# ── 自動報告生成 ─────────────────────────────────────

def report_gen(title: str, data_json: str, output: str = ""):
    try:
        import jinja2, json
        data = json.loads(data_json)
        out_path = output or str(Path.home() / "Desktop" / f"report_{dt.datetime.now().strftime('%Y%m%d_%H%M%S')}.html")
        template_str = """<!DOCTYPE html>
<html><head><meta charset="utf-8">
<style>body{font-family:sans-serif;margin:40px}table{border-collapse:collapse;width:100%}
th,td{border:1px solid #ccc;padding:8px;text-align:left}th{background:#4472C4;color:white}
tr:nth-child(even){background:#f2f2f2}h1{color:#4472C4}</style></head>
<body><h1>{{ title }}</h1><p>生成時間：{{ time }}</p>
{% for section, rows in data.items() %}
<h2>{{ section }}</h2>
{% if rows is iterable and rows is not string %}
{% if rows[0] is mapping %}
<table><tr>{% for k in rows[0].keys() %}<th>{{ k }}</th>{% endfor %}</tr>
{% for row in rows %}<tr>{% for v in row.values() %}<td>{{ v }}</td>{% endfor %}</tr>{% endfor %}
</table>
{% else %}<ul>{% for item in rows %}<li>{{ item }}</li>{% endfor %}</ul>{% endif %}
{% else %}<p>{{ rows }}</p>{% endif %}
{% endfor %}</body></html>"""
        tmpl = jinja2.Template(template_str)
        html = tmpl.render(title=title, data=data, time=dt.datetime.now().strftime("%Y-%m-%d %H:%M:%S"))
        Path(out_path).write_text(html, encoding="utf-8")
        print(f"✅ 報告已生成：{out_path}")
    except Exception as e:
        print(f"❌ 報告生成失敗：{e}")


# ── Dropbox ──────────────────────────────────────────

def dropbox_upload(local: str, remote: str, token: str = ""):
    try:
        import dropbox as dbx
        tok = token or os.getenv("DROPBOX_TOKEN", "")
        if not tok:
            print("❌ 請設定 DROPBOX_TOKEN 環境變數")
            return
        d = dbx.Dropbox(tok)
        with open(local, "rb") as f:
            d.files_upload(f.read(), remote, mode=dbx.files.WriteMode.overwrite)
        print(f"✅ 已上傳到 Dropbox：{remote}")
    except Exception as e:
        print(f"❌ Dropbox 上傳失敗：{e}")

def dropbox_download(remote: str, local: str, token: str = ""):
    try:
        import dropbox as dbx
        tok = token or os.getenv("DROPBOX_TOKEN", "")
        if not tok:
            print("❌ 請設定 DROPBOX_TOKEN 環境變數")
            return
        d = dbx.Dropbox(tok)
        _, res = d.files_download(remote)
        Path(local).write_bytes(res.content)
        print(f"✅ 已從 Dropbox 下載：{local}")
    except Exception as e:
        print(f"❌ Dropbox 下載失敗：{e}")


# ── Docker ───────────────────────────────────────────

def docker_op(action: str, name: str = ""):
    try:
        import docker as _docker
        client = _docker.from_env()
        if action == "list":
            for c in client.containers.list(all=True):
                print(f"[{c.status}] {c.name}  {c.image.tags}")
        elif action == "start":
            client.containers.get(name).start()
            print(f"✅ 容器 [{name}] 已啟動")
        elif action == "stop":
            client.containers.get(name).stop()
            print(f"✅ 容器 [{name}] 已停止")
        elif action == "logs":
            logs = client.containers.get(name).logs(tail=50).decode(errors="replace")
            print(logs)
        elif action == "images":
            for img in client.images.list():
                print(f"{img.tags}  {img.short_id}")
    except Exception as e:
        print(f"❌ Docker 操作失敗：{e}")


# ── PDF 轉圖片 ───────────────────────────────────────

def pdf_to_images(path: str, output_dir: str = "", dpi: int = 150):
    try:
        import fitz
        doc = fitz.open(path)
        out = Path(output_dir) if output_dir else Path(path).parent / (Path(path).stem + "_imgs")
        out.mkdir(parents=True, exist_ok=True)
        for i, page in enumerate(doc):
            mat = fitz.Matrix(dpi/72, dpi/72)
            pix = page.get_pixmap(matrix=mat)
            img_path = str(out / f"page_{i+1}.png")
            pix.save(img_path)
        print(f"✅ 已轉換 {len(doc)} 頁到：{out}")
    except Exception as e:
        print(f"❌ PDF 轉圖片失敗：{e}")


# ── 條碼掃描 ─────────────────────────────────────────

def barcode_scan(image_path: str = ""):
    try:
        from pyzbar.pyzbar import decode
        from PIL import Image
        img = Image.open(image_path) if image_path else pyautogui.screenshot()
        results = decode(img)
        if not results:
            print("❌ 未偵測到條碼或 QR Code")
            return
        for r in results:
            print(f"類型：{r.type}  內容：{r.data.decode('utf-8', errors='replace')}")
    except Exception as e:
        print(f"❌ 條碼掃描失敗：{e}")


# ── NLP 文字分析 ─────────────────────────────────────

def nlp_summarize(text: str):
    try:
        import anthropic
        client = anthropic.Anthropic(api_key=os.getenv("ANTHROPIC_API_KEY"))
        msg = client.messages.create(
            model="claude-haiku-4-5-20251001",
            max_tokens=512,
            messages=[{"role": "user", "content": f"請用繁體中文摘要以下文字（100字以內）：\n\n{text}"}]
        )
        print(msg.content[0].text)
    except Exception as e:
        print(f"❌ 摘要失敗：{e}")

def nlp_sentiment(text: str):
    try:
        import anthropic
        client = anthropic.Anthropic(api_key=os.getenv("ANTHROPIC_API_KEY"))
        msg = client.messages.create(
            model="claude-haiku-4-5-20251001",
            max_tokens=128,
            messages=[{"role": "user", "content": f"分析以下文字的情緒，只回覆：正面/負面/中性 + 一句說明：\n\n{text}"}]
        )
        print(msg.content[0].text)
    except Exception as e:
        print(f"❌ 情緒分析失敗：{e}")


# ── VPN 控制 ─────────────────────────────────────────

def vpn_control(action: str, name: str = "", user: str = "", password: str = ""):
    try:
        if action == "list":
            r = subprocess.run(["powershell.exe", "-Command",
                "Get-VpnConnection | Select-Object Name,ConnectionStatus | Format-Table"],
                capture_output=True, text=True, encoding="utf-8", errors="replace")
            print(r.stdout or "（未設定 VPN）")
        elif action == "connect":
            r = subprocess.run(["rasdial", name, user, password],
                capture_output=True, text=True, encoding="cp950", errors="replace")
            print(r.stdout.strip())
        elif action == "disconnect":
            r = subprocess.run(["rasdial", name, "/disconnect"],
                capture_output=True, text=True, encoding="cp950", errors="replace")
            print(r.stdout.strip())
    except Exception as e:
        print(f"❌ VPN 操作失敗：{e}")


# ── 系統還原點 ───────────────────────────────────────

def sys_restore(action: str, description: str = ""):
    try:
        if action == "create":
            ps = f"Checkpoint-Computer -Description '{description or 'Claude Auto Restore'}' -RestorePointType MODIFY_SETTINGS"
            r = subprocess.run(["powershell.exe", "-Command", ps],
                capture_output=True, text=True, encoding="utf-8", errors="replace")
            print(r.stdout or "✅ 還原點已建立")
        elif action == "list":
            r = subprocess.run(["powershell.exe", "-Command",
                "Get-ComputerRestorePoint | Select-Object SequenceNumber,Description,CreationTime | Format-Table"],
                capture_output=True, text=True, encoding="utf-8", errors="replace")
            print(r.stdout or "（無還原點）")
    except Exception as e:
        print(f"❌ 系統還原點操作失敗：{e}")


# ── 磁碟分析 ─────────────────────────────────────────

def disk_analyze(path: str = "C:/", top: int = 10):
    try:
        import psutil
        usage = psutil.disk_usage(path)
        print(f"磁碟：{path}\n總容量：{usage.total/1024**3:.1f} GB\n已使用：{usage.used/1024**3:.1f} GB ({usage.percent}%)\n可用：{usage.free/1024**3:.1f} GB\n")
        sizes = []
        try:
            for item in Path(path).iterdir():
                try:
                    if item.is_dir():
                        size = sum(f.stat().st_size for f in item.rglob("*") if f.is_file())
                    else:
                        size = item.stat().st_size
                    sizes.append((size, str(item)))
                except Exception:
                    pass
        except Exception:
            pass
        sizes.sort(reverse=True)
        print(f"佔用最多的前 {top} 個項目：")
        for size, name in sizes[:top]:
            print(f"  {size/1024**3:.2f} GB  {name}")
    except Exception as e:
        print(f"❌ 磁碟分析失敗：{e}")


# ── 人臉偵測 ─────────────────────────────────────────

def face_detect(image_path: str = "", output: str = ""):
    try:
        import cv2, numpy as np
        if not image_path:
            img_pil = pyautogui.screenshot()
            img = cv2.cvtColor(np.array(img_pil), cv2.COLOR_RGB2BGR)
        else:
            img = cv2.imread(image_path)
        cascade = cv2.CascadeClassifier(cv2.data.haarcascades + "haarcascade_frontalface_default.xml")
        gray = cv2.cvtColor(img, cv2.COLOR_BGR2GRAY)
        faces = cascade.detectMultiScale(gray, 1.1, 4)
        for (x, y, w, h) in faces:
            cv2.rectangle(img, (x, y), (x+w, y+h), (0, 255, 0), 2)
        out_path = output or str(Path.home() / "Desktop" / f"faces_{dt.datetime.now().strftime('%H%M%S')}.jpg")
        cv2.imwrite(out_path, img)
        print(f"✅ 偵測到 {len(faces)} 張人臉，已存：{out_path}")
    except Exception as e:
        print(f"❌ 人臉偵測失敗：{e}")


# ── 影片轉 GIF ───────────────────────────────────────

def video_to_gif(path: str, start: float = 0, duration: float = 5.0, output: str = "", fps: int = 10):
    try:
        import imageio
        out = output or path.replace(".mp4", ".gif")
        reader = imageio.get_reader(path)
        meta = reader.get_meta_data()
        video_fps = meta.get("fps", 30)
        start_frame = int(start * video_fps)
        end_frame = int((start + duration) * video_fps)
        frames = []
        for i, frame in enumerate(reader):
            if i < start_frame:
                continue
            if i >= end_frame:
                break
            frames.append(frame)
        imageio.mimsave(out, frames, fps=fps)
        print(f"✅ GIF 已生成：{out}（{len(frames)} 幀）")
    except Exception as e:
        print(f"❌ 影片轉 GIF 失敗：{e}")


# ── 影片生成 ─────────────────────────────────────────

def ai_video(prompt: str, provider: str = "replicate", model: str = "",
             image_url: str = "", duration: float = 5, output: str = ""):
    """
    用 AI API 生成影片
    provider: replicate / runway / kling
    需要 .env 或環境變數中對應的 API Key
    """
    import requests, time, traceback, os
    from pathlib import Path
    from dotenv import load_dotenv
    load_dotenv(Path("C:/Users/blue_/claude-telegram-bot/.env"))

    out = output or str(Path.home() / "Desktop" / f"ai_video_{int(time.time())}.mp4")

    def _download(url, dest):
        r = requests.get(url, timeout=120, stream=True)
        r.raise_for_status()
        with open(dest, "wb") as f:
            for chunk in r.iter_content(8192):
                f.write(chunk)
        return dest

    try:
        if provider == "replicate":
            api_key = os.getenv("REPLICATE_API_TOKEN", "")
            if not api_key:
                print("❌ 缺少 REPLICATE_API_TOKEN，請在 .env 加入"); return
            mdl = model or ("stability-ai/stable-video-diffusion" if image_url else "minimax/video-01")
            headers = {"Authorization": f"Bearer {api_key}", "Content-Type": "application/json"}
            inputs = {"prompt": prompt}
            if image_url: inputs["image"] = image_url
            if duration:  inputs["duration"] = int(duration)
            r = requests.post(
                f"https://api.replicate.com/v1/models/{mdl}/predictions",
                json={"input": inputs}, headers=headers, timeout=30
            )
            r.raise_for_status()
            pred_id = r.json()["id"]
            for _ in range(60):
                time.sleep(5)
                resp = requests.get(f"https://api.replicate.com/v1/predictions/{pred_id}",
                                    headers=headers, timeout=15)
                data = resp.json()
                if data.get("status") == "succeeded":
                    url = data["output"]
                    if isinstance(url, list): url = url[0]
                    _download(url, out)
                    print(f"✅ Replicate 影片已生成：{out}")
                    return
                elif data.get("status") == "failed":
                    print(f"❌ Replicate 失敗：{data.get('error')}")
                    return
            print("❌ Replicate 逾時")

        elif provider == "runway":
            api_key = os.getenv("RUNWAY_API_KEY", "")
            if not api_key:
                print("❌ 缺少 RUNWAY_API_KEY，請在 .env 加入"); return
            import runwayml
            client = runwayml.RunwayML(api_key=api_key)
            if image_url:
                task = client.image_to_video.create(
                    model="gen4_turbo", prompt_image=image_url,
                    prompt_text=prompt, duration=int(min(duration, 10)), ratio="1280:720")
            else:
                task = client.text_to_video.create(
                    model="gen4_turbo", prompt_text=prompt,
                    duration=int(min(duration, 10)), ratio="1280:720")
            for _ in range(60):
                time.sleep(5)
                t = client.tasks.retrieve(task.id)
                if t.status == "SUCCEEDED":
                    _download(t.output[0], out)
                    print(f"✅ Runway 影片已生成：{out}"); return
                elif t.status in ("FAILED", "CANCELLED"):
                    print(f"❌ Runway 失敗：{t.failure_reason}"); return
            print("❌ Runway 逾時")

        elif provider == "kling":
            import jwt as _jwt
            access_key = os.getenv("KLING_ACCESS_KEY", "")
            secret_key = os.getenv("KLING_SECRET_KEY", "")
            if not access_key or not secret_key:
                print("❌ 缺少 KLING_ACCESS_KEY / KLING_SECRET_KEY"); return
            def _token():
                return _jwt.encode({"iss": access_key, "exp": int(time.time())+1800,
                                    "nbf": int(time.time())-5}, secret_key, algorithm="HS256")
            hdrs = {"Authorization": f"Bearer {_token()}", "Content-Type": "application/json"}
            body = {"model_name": "kling-v1", "prompt": prompt,
                    "duration": str(int(min(duration, 10))), "aspect_ratio": "16:9"}
            if image_url: body["image_url"] = image_url
            ep = "image2video" if image_url else "text2video"
            r = requests.post(f"https://api.klingai.com/v1/videos/{ep}",
                              json=body, headers=hdrs, timeout=30)
            r.raise_for_status()
            d = r.json()
            if d.get("code") != 0:
                print(f"❌ Kling 錯誤：{d.get('message')}"); return
            task_id = d["data"]["task_id"]
            for _ in range(60):
                time.sleep(5)
                hdrs["Authorization"] = f"Bearer {_token()}"
                resp = requests.get(f"https://api.klingai.com/v1/videos/{ep}/{task_id}",
                                    headers=hdrs, timeout=15)
                dd = resp.json().get("data", {})
                if dd.get("task_status") == "succeed":
                    _download(dd["task_result"]["videos"][0]["url"], out)
                    print(f"✅ Kling 影片已生成：{out}"); return
                elif dd.get("task_status") == "failed":
                    print(f"❌ Kling 失敗：{dd.get('task_status_msg')}"); return
            print("❌ Kling 逾時")
        else:
            print(f"❌ 未知 provider：{provider}")
    except Exception as e:
        print(f"❌ AI 影片生成失敗：{e}\n{traceback.format_exc()}")


def video_gen(mode: str = "slideshow", output: str = "", **kwargs):
    """
    生成影片
    mode: slideshow / text_video / tts_video / screen_record
    """
    import re as _re
    import numpy as np
    import subprocess
    import tempfile
    import traceback
    from pathlib import Path
    from PIL import Image, ImageDraw, ImageFont
    import imageio_ffmpeg

    ffmpeg_exe = imageio_ffmpeg.get_ffmpeg_exe()
    w, h = kwargs.get("size", (1280, 720))
    fps  = kwargs.get("fps", 24)
    out  = output or str(Path.home() / "Desktop" / f"video_{dt.datetime.now().strftime('%H%M%S')}.mp4")

    def _write_frames(frames_iter, out_path, vid_fps, width, height):
        proc = subprocess.Popen([
            ffmpeg_exe, "-y",
            "-f", "rawvideo", "-vcodec", "rawvideo",
            "-s", f"{width}x{height}", "-pix_fmt", "rgb24",
            "-r", str(vid_fps), "-i", "pipe:0",
            "-vcodec", "libx264", "-pix_fmt", "yuv420p",
            "-preset", "fast", "-crf", "23",
            out_path
        ], stdin=subprocess.PIPE, stderr=subprocess.PIPE)
        for frame in frames_iter:
            proc.stdin.write(np.asarray(frame, dtype=np.uint8).tobytes())
        proc.stdin.close()
        rc  = proc.wait()
        err = proc.stderr.read().decode(errors="replace")
        if rc != 0 or not Path(out_path).exists():
            raise RuntimeError(f"ffmpeg 失敗 rc={rc}: {err[-300:]}")

    def _get_font(pt=60):
        for fp in ["C:/Windows/Fonts/msjh.ttc", "C:/Windows/Fonts/msyh.ttc",
                   "C:/Windows/Fonts/simhei.ttf", "C:/Windows/Fonts/arial.ttf"]:
            if Path(fp).exists():
                try:
                    return ImageFont.truetype(fp, pt)
                except Exception:
                    pass
        return ImageFont.load_default()

    def _tts_sync(text: str, voice: str, out_path: str):
        import asyncio, edge_tts
        loop = asyncio.new_event_loop()
        asyncio.set_event_loop(loop)
        try:
            loop.run_until_complete(edge_tts.Communicate(text, voice).save(out_path))
        finally:
            loop.close()
        if not Path(out_path).exists():
            raise RuntimeError("TTS 語音檔案未生成")

    def _get_audio_dur(audio_path: str) -> float:
        r = subprocess.run([ffmpeg_exe, "-i", audio_path],
                           capture_output=True, text=True, errors="replace")
        m = _re.search(r"Duration:\s*(\d+):(\d+):(\d+\.?\d*)", r.stderr)
        return int(m.group(1)) * 3600 + int(m.group(2)) * 60 + float(m.group(3)) if m else 10.0

    try:
        if mode == "slideshow":
            images = kwargs.get("images", [])
            dur    = kwargs.get("duration", 3)
            sl_fps = kwargs.get("fps", 12)
            trans  = kwargs.get("transition", 0.5)
            if not images:
                print("❌ 請提供 images 參數（圖片路徑列表）"); return

            sf = int(sl_fps * dur)
            tf = int(sl_fps * trans)

            def _gen():
                loaded = []
                for p in images:
                    try:
                        loaded.append(np.array(Image.open(p).convert("RGB").resize((w, h), Image.LANCZOS)))
                    except Exception:
                        loaded.append(np.zeros((h, w, 3), dtype=np.uint8))
                for arr in loaded:
                    for i in range(tf):  yield (arr * (i / tf)).astype(np.uint8)
                    for _ in range(max(1, sf - tf * 2)):  yield arr
                    for i in range(tf):  yield (arr * (1.0 - i / tf)).astype(np.uint8)

            _write_frames(_gen(), out, sl_fps, w, h)
            print(f"✅ 投影片影片已生成：{out}")

        elif mode == "text_video":
            text   = kwargs.get("text", "Hello")
            dur    = kwargs.get("duration", 5)
            bg_col = tuple(kwargs.get("bg_color",   [30, 30, 40]))
            fg_col = tuple(kwargs.get("font_color", [255, 255, 255]))
            fsize  = kwargs.get("font_size", 60)
            font   = _get_font(fsize)
            total  = max(1, int(fps * dur))

            def _gen():
                for i in range(total):
                    p = i / total
                    a = p / 0.2 if p < 0.2 else (1.0 - p) / 0.2 if p > 0.8 else 1.0
                    img  = Image.new("RGB", (w, h), bg_col)
                    draw = ImageDraw.Draw(img)
                    c    = tuple(int(x * a) for x in fg_col)
                    lines = text.split("\n")
                    lh    = fsize + 10
                    y0    = (h - len(lines) * lh) // 2
                    for j, ln in enumerate(lines):
                        bb = draw.textbbox((0, 0), ln, font=font)
                        draw.text(((w - (bb[2] - bb[0])) // 2, y0 + j * lh), ln, font=font, fill=c)
                    yield np.array(img)

            _write_frames(_gen(), out, fps, w, h)
            print(f"✅ 文字動畫影片已生成：{out}")

        elif mode == "tts_video":
            text     = kwargs.get("text", "")
            img_path = kwargs.get("image", "")
            voice    = kwargs.get("voice", "zh-CN-YunjianNeural")
            subtitle = kwargs.get("subtitle", True)
            if not text:
                print("❌ 請提供 text 參數"); return

            tmp    = Path(tempfile.mkdtemp())
            audio  = str(tmp / "tts.mp3")
            vidtmp = str(tmp / "silent.mp4")

            _tts_sync(clean_for_tts(text), voice, audio)
            dur   = _get_audio_dur(audio)
            total = max(1, int(fps * dur))

            if img_path and Path(img_path).exists():
                bg = np.array(Image.open(img_path).convert("RGB").resize((w, h), Image.LANCZOS))
            else:
                bg = np.zeros((h, w, 3), dtype=np.uint8)
                for row in range(h):
                    v = int(20 + 30 * row / h)
                    bg[row, :] = [v, v, v + 20]

            font = _get_font(40)
            cpl  = 20
            subs = [text[i:i+cpl] for i in range(0, len(text), cpl)]

            def _gen():
                for i in range(total):
                    frame = bg.copy()
                    if subtitle and subs:
                        ln  = subs[min(int(i / total * len(subs)), len(subs) - 1)]
                        img = Image.fromarray(frame)
                        d   = ImageDraw.Draw(img)
                        bb  = d.textbbox((0, 0), ln, font=font)
                        tw  = bb[2] - bb[0]
                        tx  = (w - tw) // 2
                        ty  = int(h * 0.82)
                        d.rectangle([tx - 10, ty - 5, tx + tw + 10, ty + 45], fill=(0, 0, 0))
                        d.text((tx, ty), ln, font=font, fill=(255, 255, 255))
                        frame = np.array(img)
                    yield frame

            _write_frames(_gen(), vidtmp, fps, w, h)
            r = subprocess.run([
                ffmpeg_exe, "-y", "-i", vidtmp, "-i", audio,
                "-c:v", "copy", "-c:a", "aac", "-shortest", out
            ], capture_output=True)
            if not Path(out).exists():
                raise RuntimeError(f"音訊合併失敗：{r.stderr.decode(errors='replace')[-200:]}")
            print(f"✅ TTS 語音影片已生成：{out}")

        elif mode == "screen_record":
            import mss, time as _time
            dur      = kwargs.get("duration", 10)
            rec_fps  = kwargs.get("fps", 10)
            interval = 1.0 / rec_fps
            total    = max(1, int(rec_fps * dur))

            with mss.mss() as sct:
                mon    = sct.monitors[1]
                sw, sh = mon["width"], mon["height"]

                def _gen():
                    for _ in range(total):
                        t0  = _time.time()
                        arr = np.array(sct.grab(mon))[:, :, :3][:, :, ::-1]
                        yield arr
                        elapsed = _time.time() - t0
                        if elapsed < interval:
                            _time.sleep(interval - elapsed)

                _write_frames(_gen(), out, rec_fps, sw, sh)
            print(f"✅ 螢幕錄影完成：{out}（{dur} 秒）")

        else:
            print(f"❌ 未知 mode：{mode}")

    except Exception as e:
        print(f"❌ 影片生成失敗：{e}\n{traceback.format_exc()}")


# ── Excel 圖表 ───────────────────────────────────────

def excel_chart(path: str, sheet: str, chart_type: str = "bar", title: str = ""):
    try:
        import openpyxl
        from openpyxl.chart import BarChart, LineChart, PieChart, Reference
        wb = openpyxl.load_workbook(path)
        ws = wb[sheet] if sheet in wb.sheetnames else wb.active
        max_row = ws.max_row
        max_col = ws.max_column
        chart_map = {"bar": BarChart, "line": LineChart, "pie": PieChart}
        chart = chart_map.get(chart_type, BarChart)()
        chart.title = title or sheet
        chart.style = 10
        data = Reference(ws, min_col=2, min_row=1, max_row=max_row, max_col=max_col)
        cats = Reference(ws, min_col=1, min_row=2, max_row=max_row)
        chart.add_data(data, titles_from_data=True)
        chart.set_categories(cats)
        ws.add_chart(chart, "A" + str(max_row + 2))
        wb.save(path)
        print(f"✅ 圖表已加入：{path}")
    except Exception as e:
        print(f"❌ Excel 圖表失敗：{e}")


# ── 網路速度測試 ─────────────────────────────────────

def speedtest_run():
    try:
        import speedtest as st
        print("⏳ 測速中（約需 30 秒）...")
        s = st.Speedtest()
        s.get_best_server()
        download = s.download() / 1_000_000
        upload = s.upload() / 1_000_000
        ping = s.results.ping
        server = s.results.server.get("name","")
        print(f"📶 網路速度測試結果\n下載：{download:.1f} Mbps\n上傳：{upload:.1f} Mbps\nPing：{ping:.0f} ms\n伺服器：{server}")
    except Exception as e:
        print(f"❌ 速度測試失敗：{e}")


# ── 螢幕截圖比對 ─────────────────────────────────────

def screenshot_compare(img1_path: str = "", img2_path: str = "", output: str = ""):
    try:
        import cv2, numpy as np
        from PIL import Image
        if not img1_path:
            img1 = np.array(pyautogui.screenshot())
        else:
            img1 = cv2.imread(img1_path)
        if not img2_path:
            import time as t; t.sleep(2)
            img2 = np.array(pyautogui.screenshot())
        else:
            img2 = cv2.imread(img2_path)
        img1_bgr = cv2.cvtColor(img1, cv2.COLOR_RGB2BGR) if img1_path == "" else img1
        img2_bgr = cv2.cvtColor(img2, cv2.COLOR_RGB2BGR) if img2_path == "" else img2
        h, w = min(img1_bgr.shape[0], img2_bgr.shape[0]), min(img1_bgr.shape[1], img2_bgr.shape[1])
        img1_bgr, img2_bgr = img1_bgr[:h,:w], img2_bgr[:h,:w]
        diff = cv2.absdiff(img1_bgr, img2_bgr)
        gray = cv2.cvtColor(diff, cv2.COLOR_BGR2GRAY)
        _, thresh = cv2.threshold(gray, 30, 255, cv2.THRESH_BINARY)
        contours, _ = cv2.findContours(thresh, cv2.RETR_EXTERNAL, cv2.CHAIN_APPROX_SIMPLE)
        result = img2_bgr.copy()
        cv2.drawContours(result, contours, -1, (0, 0, 255), 2)
        out = output or str(Path.home() / "Desktop" / f"diff_{dt.datetime.now().strftime('%H%M%S')}.png")
        cv2.imwrite(out, result)
        changed = cv2.countNonZero(thresh)
        total = h * w
        pct = changed / total * 100
        print(f"✅ 差異：{pct:.2f}%，已標記差異區域：{out}")
    except Exception as e:
        print(f"❌ 截圖比對失敗：{e}")


# ── 一次性提醒 ───────────────────────────────────────

def set_reminder(time_str: str, message: str):
    """time_str: HH:MM 或 秒數（如 '300' 代表5分鐘後）"""
    try:
        import threading, time as t
        def _remind():
            if time_str.isdigit():
                t.sleep(int(time_str))
            else:
                import datetime as dt
                now = dt.dt.datetime.now()
                target = dt.dt.datetime.strptime(time_str, "%H:%M").replace(
                    year=now.year, month=now.month, day=now.day)
                if target < now:
                    target = target.replace(day=now.day + 1)
                t.sleep((target - now).total_seconds())
            try:
                from win10toast import ToastNotifier
                ToastNotifier().show_toast("⏰ 提醒", message, duration=10)
            except Exception:
                pass
            try:
                import pyttsx3
                engine = pyttsx3.init()
                engine.say(message)
                engine.runAndWait()
            except Exception:
                pass
            print(f"⏰ 提醒：{message}")
        threading.Thread(target=_remind, daemon=True).start()
        print(f"✅ 提醒已設定：{time_str} → {message}")
    except Exception as e:
        print(f"❌ 設定提醒失敗：{e}")


# ── 網頁全頁截圖 ─────────────────────────────────────

def webpage_screenshot(url: str, output: str = ""):
    try:
        from playwright.sync_api import sync_playwright
        out = output or str(Path.home() / "Desktop" / f"webpage_{dt.datetime.now().strftime('%H%M%S')}.png")
        with sync_playwright() as p:
            browser = p.chromium.launch()
            page = browser.new_page(viewport={"width": 1280, "height": 800})
            page.goto(url, timeout=30000)
            page.wait_for_load_state("networkidle")
            page.screenshot(path=out, full_page=True)
            browser.close()
        print(f"✅ 網頁截圖已存：{out}")
    except Exception as e:
        print(f"❌ 網頁截圖失敗：{e}")


# ── 網頁變化監控 ─────────────────────────────────────

def web_monitor(url: str, selector: str = "body", interval: float = 60.0, duration: float = 3600.0):
    try:
        import hashlib, time as t
        from bs4 import BeautifulSoup
        def _fetch():
            r = requests.get(url, timeout=15, headers={"User-Agent": "Mozilla/5.0"})
            soup = BeautifulSoup(r.text, "html.parser")
            elements = soup.select(selector)
            text = "\n".join(e.get_text(strip=True) for e in elements)
            return hashlib.md5(text.encode()).hexdigest(), text[:200]
        last_hash, _ = _fetch()
        print(f"🔍 開始監控：{url}（每 {interval} 秒，持續 {duration} 秒）")
        end = t.time() + duration
        changes = 0
        while t.time() < end:
            t.sleep(interval)
            try:
                new_hash, snippet = _fetch()
                if new_hash != last_hash:
                    changes += 1
                    print(f"⚠️ [{dt.datetime.now().strftime('%H:%M:%S')}] 網頁有變化！\n{snippet}")
                    last_hash = new_hash
            except Exception as e:
                print(f"檢查失敗：{e}")
        print(f"✅ 監控結束，共偵測到 {changes} 次變化")
    except Exception as e:
        print(f"❌ 網頁監控失敗：{e}")


# ── 批次重新命名 ─────────────────────────────────────

def batch_rename(folder: str, pattern: str, replacement: str, ext_filter: str = ""):
    try:
        import re
        folder_path = Path(folder)
        files = [f for f in folder_path.iterdir() if f.is_file()]
        if ext_filter:
            files = [f for f in files if f.suffix.lower() == ext_filter.lower()]
        count = 0
        for f in sorted(files):
            new_name = re.sub(pattern, replacement, f.stem) + f.suffix
            new_path = f.parent / new_name
            if new_path != f:
                f.rename(new_path)
                print(f"  {f.name} → {new_name}")
                count += 1
        print(f"✅ 已重新命名 {count} 個檔案")
    except Exception as e:
        print(f"❌ 批次重新命名失敗：{e}")


# ── 圖片壓縮 ─────────────────────────────────────────

def img_compress(path: str, quality: int = 75, output: str = ""):
    try:
        from PIL import Image
        img = Image.open(path)
        out = output or path.replace(".", f"_q{quality}.")
        if img.mode in ("RGBA", "P"):
            img = img.convert("RGB")
        img.save(out, optimize=True, quality=quality)
        orig_size = Path(path).stat().st_size
        new_size = Path(out).stat().st_size
        print(f"✅ 壓縮完成：{orig_size//1024}KB → {new_size//1024}KB（{(1-new_size/orig_size)*100:.1f}% 節省）\n儲存至：{out}")
    except Exception as e:
        print(f"❌ 圖片壓縮失敗：{e}")

def batch_img_process(folder: str, action: str, width: int = 0, height: int = 0, quality: int = 75):
    try:
        from PIL import Image
        folder_path = Path(folder)
        out_dir = folder_path / f"output_{action}"
        out_dir.mkdir(exist_ok=True)
        count = 0
        for f in folder_path.glob("*.{jpg,jpeg,png,bmp,webp}"):
            try:
                img = Image.open(f)
                if action == "resize" and width and height:
                    img = img.resize((width, height))
                elif action == "compress":
                    if img.mode in ("RGBA", "P"): img = img.convert("RGB")
                img.save(str(out_dir / f.name), optimize=True, quality=quality)
                count += 1
            except Exception:
                pass
        print(f"✅ 批次處理完成：{count} 張圖片 → {out_dir}")
    except Exception as e:
        print(f"❌ 批次圖片處理失敗：{e}")


# ── OCR + 翻譯 pipeline ──────────────────────────────

def ocr_translate(image_path: str = "", target: str = "zh-TW"):
    try:
        import easyocr, numpy as np
        from PIL import Image
        from deep_translator import GoogleTranslator
        reader = easyocr.Reader(["ch_tra", "en"], gpu=False)
        img = Image.open(image_path) if image_path else pyautogui.screenshot()
        results = reader.readtext(np.array(img))
        text = " ".join([r[1] for r in results])
        if not text.strip():
            print("❌ 未辨識到文字")
            return
        print(f"原文：{text[:300]}")
        translated = GoogleTranslator(source="auto", target=target).translate(text)
        print(f"翻譯（{target}）：{translated}")
    except Exception as e:
        print(f"❌ OCR 翻譯失敗：{e}")


# ── IP 資訊查詢 ──────────────────────────────────────

def ip_info(ip: str = ""):
    try:
        target = ip or ""
        res = requests.get(f"http://ip-api.com/json/{target}?lang=zh-TW", timeout=10)
        d = res.json()
        if d.get("status") == "fail":
            print(f"❌ 查詢失敗：{d.get('message')}")
            return
        print(
            f"IP：{d.get('query')}\n"
            f"國家：{d.get('country')} ({d.get('countryCode')})\n"
            f"城市：{d.get('city')}  地區：{d.get('regionName')}\n"
            f"ISP：{d.get('isp')}\n"
            f"組織：{d.get('org')}\n"
            f"時區：{d.get('timezone')}\n"
            f"座標：{d.get('lat')}, {d.get('lon')}"
        )
    except Exception as e:
        print(f"❌ IP 查詢失敗：{e}")


# ── 匯率查詢 ─────────────────────────────────────────

def currency(amount: float, from_cur: str, to_cur: str):
    try:
        res = requests.get(
            f"https://api.frankfurter.app/latest?amount={amount}&from={from_cur.upper()}&to={to_cur.upper()}",
            timeout=10
        )
        data = res.json()
        rate = data["rates"].get(to_cur.upper())
        if rate is None:
            print(f"❌ 找不到 {to_cur} 匯率")
            return
        print(f"💱 {amount} {from_cur.upper()} = {rate:.4f} {to_cur.upper()}")
        print(f"（匯率基準日：{data.get('date')}）")
    except Exception as e:
        print(f"❌ 匯率查詢失敗：{e}")


# ── Windows 事件日誌 ─────────────────────────────────

def event_log(log_name: str = "System", level: str = "Error", count: int = 10):
    try:
        ps = (
            f"Get-WinEvent -LogName '{log_name}' -MaxEvents {count} "
            f"| Where-Object {{$_.LevelDisplayName -eq '{level}'}} "
            f"| Select-Object TimeCreated,Id,Message "
            f"| Format-List"
        )
        r = subprocess.run(["powershell.exe", "-Command", ps],
            capture_output=True, text=True, encoding="utf-8", errors="replace", timeout=30)
        print(r.stdout[:3000] or f"（無 {level} 等級事件）")
    except Exception as e:
        print(f"❌ 事件日誌查詢失敗：{e}")


# ── 進階 TTS (edge-tts) ──────────────────────────────

def tts_edge(text: str, voice: str = "zh-TW-HsiaoChenNeural", output: str = ""):
    try:
        import edge_tts, asyncio
        out = output or str(Path.home() / "Desktop" / f"tts_{dt.datetime.now().strftime('%H%M%S')}.mp3")
        _clean = clean_for_tts(text)
        async def _gen():
            communicate = edge_tts.Communicate(_clean, voice)
            await communicate.save(out)
        asyncio.run(_gen())
        subprocess.Popen(["powershell.exe", "-Command", f"Start-Process '{out}'"])
        print(f"✅ 語音已生成並播放：{out}")
    except Exception as e:
        print(f"❌ Edge TTS 失敗：{e}")

def tts_voices():
    try:
        import edge_tts, asyncio
        async def _list():
            return await edge_tts.list_voices()
        voices = asyncio.run(_list())
        zh_voices = [v for v in voices if v["Locale"].startswith("zh")]
        for v in zh_voices:
            print(f"{v['ShortName']}  {v['FriendlyName']}")
    except Exception as e:
        print(f"❌ 列出語音失敗：{e}")


# ── 郵件附件發送 ─────────────────────────────────────

def send_email_attach(to: str, subject: str, body: str, attachments: str = ""):
    try:
        import smtplib
        from email.mime.multipart import MIMEMultipart
        from email.mime.text import MIMEText
        from email.mime.base import MIMEBase
        from email import encoders
        smtp_user = os.getenv("SMTP_USER", "")
        smtp_pass = os.getenv("SMTP_PASS", "")
        smtp_host = os.getenv("SMTP_HOST", "smtp.gmail.com")
        smtp_port = int(os.getenv("SMTP_PORT", "587"))
        msg = MIMEMultipart()
        msg["From"] = smtp_user
        msg["To"] = to
        msg["Subject"] = subject
        msg.attach(MIMEText(body, "plain", "utf-8"))
        if attachments:
            for path in attachments.split(","):
                path = path.strip()
                if Path(path).exists():
                    with open(path, "rb") as f:
                        part = MIMEBase("application", "octet-stream")
                        part.set_payload(f.read())
                    encoders.encode_base64(part)
                    part.add_header("Content-Disposition", f"attachment; filename={Path(path).name}")
                    msg.attach(part)
        with smtplib.SMTP(smtp_host, smtp_port) as s:
            s.starttls()
            s.login(smtp_user, smtp_pass)
            s.send_message(msg)
        print(f"✅ 郵件已發送至 {to}（附件：{attachments or '無'}）")
    except Exception as e:
        print(f"❌ 發送失敗：{e}")


# ── 剪貼簿圖片 ───────────────────────────────────────

def clipboard_img_get(output: str = ""):
    try:
        import win32clipboard
        from PIL import Image
        import io as _io
        win32clipboard.OpenClipboard()
        try:
            data = win32clipboard.GetClipboardData(win32clipboard.CF_DIB)
        finally:
            win32clipboard.CloseClipboard()
        out = output or str(Path.home() / "Desktop" / f"clipboard_{dt.datetime.now().strftime('%H%M%S')}.png")
        img = Image.open(_io.BytesIO(data))
        img.save(out)
        print(f"✅ 剪貼簿圖片已存：{out}")
    except Exception as e:
        print(f"❌ 讀取剪貼簿圖片失敗：{e}")

def clipboard_img_set(path: str):
    try:
        import win32clipboard
        from PIL import Image
        import io as _io
        img = Image.open(path).convert("RGB")
        buf = _io.BytesIO()
        img.save(buf, "BMP")
        data = buf.getvalue()[14:]
        win32clipboard.OpenClipboard()
        try:
            win32clipboard.EmptyClipboard()
            win32clipboard.SetClipboardData(win32clipboard.CF_DIB, data)
        finally:
            win32clipboard.CloseClipboard()
        print(f"✅ 圖片已複製到剪貼簿：{path}")
    except Exception as e:
        print(f"❌ 設定剪貼簿圖片失敗：{e}")


# ── USB 裝置管理 ─────────────────────────────────────

def usb_list():
    try:
        r = subprocess.run(["powershell.exe", "-Command",
            "Get-PnpDevice | Where-Object {$_.Class -eq 'USB' -and $_.Status -eq 'OK'} | Select-Object FriendlyName,InstanceId | Format-Table -AutoSize"],
            capture_output=True, text=True, encoding="utf-8", errors="replace")
        print(r.stdout[:3000] or "（無 USB 裝置）")
    except Exception as e:
        print(f"❌ USB 查詢失敗：{e}")


# ── Windows 防火牆 ───────────────────────────────────

def firewall(action: str, name: str = "", direction: str = "in", port: int = 0, protocol: str = "TCP"):
    try:
        if action == "list":
            r = subprocess.run(["powershell.exe", "-Command",
                "Get-NetFirewallRule | Where-Object {$_.Enabled -eq 'True'} | Select-Object DisplayName,Direction,Action | Format-Table -AutoSize"],
                capture_output=True, text=True, encoding="utf-8", errors="replace")
            print(r.stdout[:3000])
        elif action == "add":
            ps = f"New-NetFirewallRule -DisplayName '{name}' -Direction {direction} -Protocol {protocol} -LocalPort {port} -Action Allow"
            r = subprocess.run(["powershell.exe", "-Command", ps], capture_output=True, text=True, encoding="utf-8", errors="replace")
            print(r.stdout or f"✅ 防火牆規則已新增：{name}")
        elif action == "remove":
            r = subprocess.run(["powershell.exe", "-Command", f"Remove-NetFirewallRule -DisplayName '{name}'"],
                capture_output=True, text=True, encoding="utf-8", errors="replace")
            print(r.stdout or f"✅ 防火牆規則已刪除：{name}")
    except Exception as e:
        print(f"❌ 防火牆操作失敗：{e}")


# ── 任務清單 ─────────────────────────────────────────

TODO_DB = Path.home() / "claude-telegram-bot" / "todo.db"

def _todo_init():
    conn = sqlite3.connect(str(TODO_DB))
    conn.execute("CREATE TABLE IF NOT EXISTS todos (id INTEGER PRIMARY KEY AUTOINCREMENT, task TEXT, done INTEGER DEFAULT 0, created TEXT)")
    conn.commit(); conn.close()

def todo(action: str, task: str = "", todo_id: int = 0):
    try:
        _todo_init()
        conn = sqlite3.connect(str(TODO_DB))
        cur = conn.cursor()
        if action == "add":
            cur.execute("INSERT INTO todos (task, created) VALUES (?, ?)", (task, dt.datetime.now().strftime("%Y-%m-%d %H:%M")))
            conn.commit()
            print(f"✅ 已新增任務：{task}")
        elif action == "list":
            rows = cur.execute("SELECT id, task, done, created FROM todos ORDER BY done, id").fetchall()
            if not rows: print("（清單為空）"); conn.close(); return
            for r in rows:
                status = "✅" if r[2] else "⬜"
                print(f"{status} [{r[0]}] {r[1]}  ({r[3]})")
        elif action == "done":
            cur.execute("UPDATE todos SET done=1 WHERE id=?", (todo_id,))
            conn.commit(); print(f"✅ 任務 #{todo_id} 已完成")
        elif action == "delete":
            cur.execute("DELETE FROM todos WHERE id=?", (todo_id,))
            conn.commit(); print(f"✅ 任務 #{todo_id} 已刪除")
        elif action == "clear":
            cur.execute("DELETE FROM todos WHERE done=1")
            conn.commit(); print("✅ 已清除所有已完成任務")
        conn.close()
    except Exception as e:
        print(f"❌ 任務清單操作失敗：{e}")


# ── 檔案同步 ─────────────────────────────────────────

def file_sync(src: str, dest: str, dry_run: bool = False):
    try:
        import filecmp, shutil
        src_path = Path(src); dest_path = Path(dest)
        dest_path.mkdir(parents=True, exist_ok=True)
        copied = deleted = 0
        for item in src_path.rglob("*"):
            rel = item.relative_to(src_path)
            dst = dest_path / rel
            if item.is_dir():
                dst.mkdir(parents=True, exist_ok=True)
            elif item.is_file():
                if not dst.exists() or not filecmp.cmp(str(item), str(dst), shallow=False):
                    if not dry_run:
                        shutil.copy2(str(item), str(dst))
                    print(f"  複製：{rel}")
                    copied += 1
        print(f"{'[預覽] ' if dry_run else ''}✅ 同步完成：{copied} 個檔案更新，{deleted} 個刪除")
    except Exception as e:
        print(f"❌ 檔案同步失敗：{e}")


# ── 系統資源圖表 ─────────────────────────────────────

def sysres_chart(duration: int = 10, output: str = ""):
    try:
        import psutil, matplotlib
        matplotlib.use("Agg")
        import matplotlib.pyplot as plt, time as t
        cpu_vals, mem_vals, times = [], [], []
        for i in range(duration):
            cpu_vals.append(psutil.cpu_percent(interval=1))
            mem_vals.append(psutil.virtual_memory().percent)
            times.append(i + 1)
        out = output or str(Path.home() / "Desktop" / f"sysres_{dt.datetime.now().strftime('%H%M%S')}.png")
        fig, ax = plt.subplots()
        ax.plot(times, cpu_vals, label="CPU %", color="blue")
        ax.plot(times, mem_vals, label="RAM %", color="orange")
        ax.set_ylim(0, 100); ax.set_xlabel("秒"); ax.set_ylabel("%")
        ax.set_title("系統資源使用率"); ax.legend(); plt.tight_layout()
        plt.savefig(out); plt.close()
        print(f"✅ 資源圖表已存：{out}")
    except Exception as e:
        print(f"❌ 資源圖表失敗：{e}")


# ── 密碼管理 ─────────────────────────────────────────

PWD_DB = Path.home() / "claude-telegram-bot" / "passwords.db"

def _pwd_get_key(master: str) -> bytes:
    from cryptography.hazmat.primitives.kdf.pbkdf2 import PBKDF2HMAC
    from cryptography.hazmat.primitives import hashes
    import base64
    kdf = PBKDF2HMAC(algorithm=hashes.SHA256(), length=32, salt=b"pwd_manager_v1", iterations=100000)
    return base64.urlsafe_b64encode(kdf.derive(master.encode()))

def password_save(site: str, username: str, password: str, master: str):
    try:
        from cryptography.fernet import Fernet
        key = _pwd_get_key(master)
        f = Fernet(key)
        encrypted = f.encrypt(password.encode()).decode()
        conn = sqlite3.connect(str(PWD_DB))
        conn.execute("CREATE TABLE IF NOT EXISTS passwords (id INTEGER PRIMARY KEY AUTOINCREMENT, site TEXT, username TEXT, password TEXT)")
        conn.execute("INSERT OR REPLACE INTO passwords (site, username, password) VALUES (?, ?, ?)", (site, username, encrypted))
        conn.commit(); conn.close()
        print(f"✅ 密碼已加密儲存：{site} / {username}")
    except Exception as e:
        print(f"❌ 儲存密碼失敗：{e}")

def password_get(site: str, master: str):
    try:
        from cryptography.fernet import Fernet
        key = _pwd_get_key(master)
        f = Fernet(key)
        conn = sqlite3.connect(str(PWD_DB))
        rows = conn.execute("SELECT site, username, password FROM passwords WHERE site LIKE ?", (f"%{site}%",)).fetchall()
        conn.close()
        if not rows: print(f"（找不到 {site} 的密碼）"); return
        for site_, user, enc_pwd in rows:
            try:
                pwd = f.decrypt(enc_pwd.encode()).decode()
                print(f"🔑 {site_}\n帳號：{user}\n密碼：{pwd}")
            except Exception:
                print(f"❌ 解密失敗（主密碼錯誤？）")
    except Exception as e:
        print(f"❌ 讀取密碼失敗：{e}")


# ── RDP 遠端桌面 ─────────────────────────────────────

def rdp_connect(host: str, user: str = "", width: int = 1280, height: int = 720):
    try:
        args = ["/v:" + host, f"/w:{width}", f"/h:{height}"]
        if user:
            args.append(f"/u:{user}")
        subprocess.Popen(["mstsc"] + args)
        print(f"✅ 正在連線 RDP：{host}")
    except Exception as e:
        print(f"❌ RDP 連線失敗：{e}")


# ── Chrome 書籤 ──────────────────────────────────────

def chrome_bookmarks():
    try:
        import json
        bookmark_path = Path.home() / "AppData/Local/Google/Chrome/User Data/Default/Bookmarks"
        if not bookmark_path.exists():
            print("❌ 找不到 Chrome 書籤檔案")
            return
        data = json.loads(bookmark_path.read_text(encoding="utf-8"))
        def _print_node(node, indent=0):
            if node.get("type") == "url":
                print("  " * indent + f"🔗 {node['name']}  {node['url']}")
            elif node.get("type") == "folder":
                print("  " * indent + f"📁 {node['name']}")
                for child in node.get("children", []):
                    _print_node(child, indent + 1)
        for root in data["roots"].values():
            _print_node(root)
    except Exception as e:
        print(f"❌ 讀取書籤失敗：{e}")


# ── 印表機管理 ───────────────────────────────────────

def printer_list():
    try:
        r = subprocess.run(["powershell.exe", "-Command",
            "Get-Printer | Select-Object Name,DriverName,PrinterStatus | Format-Table -AutoSize"],
            capture_output=True, text=True, encoding="utf-8", errors="replace")
        print(r.stdout or "（無印表機）")
    except Exception as e:
        print(f"❌ 印表機查詢失敗：{e}")

def printer_jobs():
    try:
        r = subprocess.run(["powershell.exe", "-Command",
            "Get-PrintJob -PrinterName (Get-Printer | Select-Object -First 1 -ExpandProperty Name) | Format-Table"],
            capture_output=True, text=True, encoding="utf-8", errors="replace")
        print(r.stdout or "（列印佇列為空）")
    except Exception as e:
        print(f"❌ 列印佇列查詢失敗：{e}")


# ── 網路芳鄰 ─────────────────────────────────────────

def net_share(action: str, share_path: str = "", drive: str = "Z:", user: str = "", password: str = ""):
    try:
        if action == "list":
            r = subprocess.run(["net", "use"], capture_output=True, text=True, encoding="cp950", errors="replace")
            print(r.stdout)
        elif action == "connect":
            args = ["net", "use", drive, share_path]
            if user: args += [f"/user:{user}", password]
            r = subprocess.run(args, capture_output=True, text=True, encoding="cp950", errors="replace")
            print(r.stdout or f"✅ 已連線 {share_path} → {drive}")
        elif action == "disconnect":
            r = subprocess.run(["net", "use", drive, "/delete"], capture_output=True, text=True, encoding="cp950", errors="replace")
            print(r.stdout or f"✅ 已中斷 {drive}")
    except Exception as e:
        print(f"❌ 網路芳鄰操作失敗：{e}")


# ── 字型列表 ─────────────────────────────────────────

def font_list(keyword: str = ""):
    try:
        r = subprocess.run(["powershell.exe", "-Command",
            "[System.Reflection.Assembly]::LoadWithPartialName('System.Drawing') | Out-Null; "
            "[System.Drawing.FontFamily]::Families | Select-Object -ExpandProperty Name"],
            capture_output=True, text=True, encoding="utf-8", errors="replace")
        fonts = [f for f in r.stdout.strip().splitlines() if keyword.lower() in f.lower()] if keyword else r.stdout.strip().splitlines()
        print("\n".join(fonts[:50]))
        if len(fonts) > 50:
            print(f"...（共 {len(fonts)} 個字型）")
    except Exception as e:
        print(f"❌ 字型列表失敗：{e}")


# ── 音量控制 ─────────────────────────────────────────

def volume_get():
    try:
        from ctypes import cast, POINTER
        from comtypes import CLSCTX_ALL
        from pycaw.pycaw import AudioUtilities, IAudioEndpointVolume
        devices = AudioUtilities.GetSpeakers()
        interface = devices.Activate(IAudioEndpointVolume._iid_, CLSCTX_ALL, None)
        volume = cast(interface, POINTER(IAudioEndpointVolume))
        muted = volume.GetMute()
        level = round(volume.GetMasterVolumeLevelScalar() * 100)
        print(f"🔊 音量：{level}%  {'（靜音）' if muted else ''}")
    except Exception as e:
        print(f"❌ 音量查詢失敗：{e}")

def volume_set(level: int):
    try:
        from ctypes import cast, POINTER
        from comtypes import CLSCTX_ALL
        from pycaw.pycaw import AudioUtilities, IAudioEndpointVolume
        devices = AudioUtilities.GetSpeakers()
        interface = devices.Activate(IAudioEndpointVolume._iid_, CLSCTX_ALL, None)
        volume = cast(interface, POINTER(IAudioEndpointVolume))
        volume.SetMasterVolumeLevelScalar(max(0.0, min(1.0, level / 100)), None)
        print(f"✅ 音量已設定為 {level}%")
    except Exception as e:
        print(f"❌ 設定音量失敗：{e}")

def volume_mute(mute: bool = True):
    try:
        from ctypes import cast, POINTER
        from comtypes import CLSCTX_ALL
        from pycaw.pycaw import AudioUtilities, IAudioEndpointVolume
        devices = AudioUtilities.GetSpeakers()
        interface = devices.Activate(IAudioEndpointVolume._iid_, CLSCTX_ALL, None)
        volume = cast(interface, POINTER(IAudioEndpointVolume))
        volume.SetMute(1 if mute else 0, None)
        print(f"✅ {'靜音' if mute else '取消靜音'}")
    except Exception as e:
        print(f"❌ 靜音操作失敗：{e}")


# ── 螢幕亮度/解析度 ──────────────────────────────────

def brightness_get():
    try:
        import screen_brightness_control as sbc
        b = sbc.get_brightness()
        print(f"💡 亮度：{b}%")
    except Exception as e:
        print(f"❌ 亮度查詢失敗：{e}")

def brightness_set(level: int):
    try:
        import screen_brightness_control as sbc
        sbc.set_brightness(max(0, min(100, level)))
        print(f"✅ 亮度已設定為 {level}%")
    except Exception as e:
        print(f"❌ 設定亮度失敗：{e}")

def resolution_list():
    try:
        r = subprocess.run(["powershell.exe", "-Command",
            "Get-CimInstance -ClassName Win32_VideoController | Select-Object CurrentHorizontalResolution,CurrentVerticalResolution,CurrentRefreshRate | Format-List"],
            capture_output=True, text=True, encoding="utf-8", errors="replace")
        print(r.stdout.strip())
    except Exception as e:
        print(f"❌ 解析度查詢失敗：{e}")


# ── 媒體播放控制 ─────────────────────────────────────

def media_control(action: str):
    """action: play_pause / next / prev / stop / volume_up / volume_down / mute"""
    try:
        import keyboard as kb
        key_map = {
            "play_pause": "play/pause media",
            "next": "next track",
            "prev": "previous track",
            "stop": "stop media",
            "volume_up": "volume up",
            "volume_down": "volume down",
            "mute": "volume mute",
        }
        key = key_map.get(action)
        if not key:
            print(f"❌ 未知動作：{action}，可用：{list(key_map.keys())}")
            return
        kb.send(key)
        print(f"✅ 媒體控制：{action}")
    except Exception as e:
        print(f"❌ 媒體控制失敗：{e}")


# ── 音訊裝置切換 ─────────────────────────────────────

def audio_devices():
    try:
        from pycaw.pycaw import AudioUtilities
        devices = AudioUtilities.GetAllDevices()
        for i, d in enumerate(devices):
            print(f"[{i}] {d.FriendlyName}")
    except Exception as e:
        print(f"❌ 音訊裝置查詢失敗：{e}")

def audio_switch(device_name: str):
    try:
        r = subprocess.run(["powershell.exe", "-Command",
            f"Get-AudioDevice -List | Where-Object {{$_.Name -like '*{device_name}*'}} | Set-AudioDevice"],
            capture_output=True, text=True, encoding="utf-8", errors="replace")
        print(r.stdout or f"✅ 已切換至：{device_name}")
    except Exception as e:
        print(f"❌ 切換音訊裝置失敗（需安裝 AudioDeviceCmdlets）：{e}")


# ── 已安裝軟體管理 ───────────────────────────────────

def software_list(keyword: str = ""):
    try:
        r = subprocess.run(["powershell.exe", "-Command",
            "Get-Package | Select-Object Name,Version | Sort-Object Name | Format-Table -AutoSize"],
            capture_output=True, text=True, encoding="utf-8", errors="replace", timeout=30)
        lines = r.stdout.strip().splitlines()
        if keyword:
            lines = [l for l in lines if keyword.lower() in l.lower()]
        print("\n".join(lines[:50]))
        if len(lines) > 50:
            print(f"...（共 {len(lines)} 個）")
    except Exception as e:
        print(f"❌ 軟體列表查詢失敗：{e}")

def software_install(name: str):
    try:
        r = subprocess.run(["powershell.exe", "-Command", f"winget install --id '{name}' --silent --accept-source-agreements --accept-package-agreements"],
            capture_output=True, text=True, encoding="utf-8", errors="replace", timeout=300)
        print(r.stdout or r.stderr or f"安裝完成：{name}")
    except Exception as e:
        print(f"❌ 安裝失敗：{e}")

def software_uninstall(name: str):
    try:
        r = subprocess.run(["powershell.exe", "-Command", f"winget uninstall --name '{name}' --silent"],
            capture_output=True, text=True, encoding="utf-8", errors="replace", timeout=300)
        print(r.stdout or r.stderr or f"✅ 已卸載：{name}")
    except Exception as e:
        print(f"❌ 卸載失敗：{e}")


# ── 開機自啟動管理 ───────────────────────────────────

def startup_list():
    try:
        r = subprocess.run(["powershell.exe", "-Command",
            "Get-CimInstance Win32_StartupCommand | Select-Object Name,Command,Location | Format-Table -AutoSize"],
            capture_output=True, text=True, encoding="utf-8", errors="replace")
        print(r.stdout.strip() or "（無開機自啟動程式）")
    except Exception as e:
        print(f"❌ 查詢失敗：{e}")

def startup_add(name: str, command: str):
    try:
        import winreg
        key = winreg.OpenKey(winreg.HKEY_CURRENT_USER,
            r"Software\Microsoft\Windows\CurrentVersion\Run", 0, winreg.KEY_SET_VALUE)
        winreg.SetValueEx(key, name, 0, winreg.REG_SZ, command)
        winreg.CloseKey(key)
        print(f"✅ 已新增開機自啟動：{name}")
    except Exception as e:
        print(f"❌ 新增失敗：{e}")

def startup_remove(name: str):
    try:
        import winreg
        key = winreg.OpenKey(winreg.HKEY_CURRENT_USER,
            r"Software\Microsoft\Windows\CurrentVersion\Run", 0, winreg.KEY_SET_VALUE)
        winreg.DeleteValue(key, name)
        winreg.CloseKey(key)
        print(f"✅ 已移除開機自啟動：{name}")
    except Exception as e:
        print(f"❌ 移除失敗：{e}")


# ── 環境變數管理 ─────────────────────────────────────

def env_get(name: str = ""):
    try:
        if name:
            val = os.environ.get(name, "（未設定）")
            print(f"{name} = {val}")
        else:
            for k, v in sorted(os.environ.items()):
                print(f"{k} = {v[:80]}")
    except Exception as e:
        print(f"❌ 環境變數查詢失敗：{e}")

def env_set(name: str, value: str, permanent: bool = False):
    try:
        os.environ[name] = value
        if permanent:
            r = subprocess.run(["powershell.exe", "-Command",
                f"[System.Environment]::SetEnvironmentVariable('{name}','{value}','User')"],
                capture_output=True, text=True, encoding="utf-8", errors="replace")
            print(f"✅ 環境變數已永久設定：{name} = {value}")
        else:
            print(f"✅ 環境變數已設定（本次）：{name} = {value}")
    except Exception as e:
        print(f"❌ 設定環境變數失敗：{e}")


# ── 使用者帳戶管理 ───────────────────────────────────

def user_list():
    try:
        r = subprocess.run(["powershell.exe", "-Command",
            "Get-LocalUser | Select-Object Name,Enabled,LastLogon | Format-Table -AutoSize"],
            capture_output=True, text=True, encoding="utf-8", errors="replace")
        print(r.stdout.strip())
    except Exception as e:
        print(f"❌ 查詢失敗：{e}")

def user_create(username: str, password: str):
    try:
        ps = f"New-LocalUser -Name '{username}' -Password (ConvertTo-SecureString '{password}' -AsPlainText -Force) -FullName '{username}'"
        r = subprocess.run(["powershell.exe", "-Command", ps],
            capture_output=True, text=True, encoding="utf-8", errors="replace")
        print(r.stdout or f"✅ 使用者已建立：{username}")
    except Exception as e:
        print(f"❌ 建立使用者失敗：{e}")

def user_delete(username: str):
    try:
        r = subprocess.run(["powershell.exe", "-Command", f"Remove-LocalUser -Name '{username}'"],
            capture_output=True, text=True, encoding="utf-8", errors="replace")
        print(r.stdout or f"✅ 使用者已刪除：{username}")
    except Exception as e:
        print(f"❌ 刪除使用者失敗：{e}")


# ── Windows 更新 ─────────────────────────────────────

def win_update(action: str = "list"):
    try:
        if action == "list":
            r = subprocess.run(["powershell.exe", "-Command",
                "Get-WindowsUpdate -AcceptAll 2>$null | Select-Object Title,Size | Format-Table -AutoSize"],
                capture_output=True, text=True, encoding="utf-8", errors="replace", timeout=60)
            print(r.stdout or "（需安裝 PSWindowsUpdate 模組）\n執行：Install-Module PSWindowsUpdate -Force")
        elif action == "install":
            r = subprocess.run(["powershell.exe", "-Command",
                "Install-WindowsUpdate -AcceptAll -AutoReboot:$false 2>&1"],
                capture_output=True, text=True, encoding="utf-8", errors="replace", timeout=600)
            print(r.stdout or "✅ 更新完成")
        elif action == "check":
            r = subprocess.run(["powershell.exe", "-Command",
                "(New-Object -ComObject Microsoft.Update.AutoUpdate).DetectNow()"],
                capture_output=True, text=True, encoding="utf-8", errors="replace")
            print("✅ 已觸發 Windows Update 檢查")
    except Exception as e:
        print(f"❌ Windows 更新失敗：{e}")


# ── 裝置管理員 ───────────────────────────────────────

def device_list(keyword: str = ""):
    try:
        ps = "Get-PnpDevice | Select-Object Status,Class,FriendlyName | Format-Table -AutoSize"
        r = subprocess.run(["powershell.exe", "-Command", ps],
            capture_output=True, text=True, encoding="utf-8", errors="replace")
        lines = r.stdout.strip().splitlines()
        if keyword:
            lines = [l for l in lines if keyword.lower() in l.lower()]
        print("\n".join(lines[:50]))
    except Exception as e:
        print(f"❌ 裝置查詢失敗：{e}")

def device_toggle(name: str, enable: bool = True):
    try:
        action = "Enable" if enable else "Disable"
        r = subprocess.run(["powershell.exe", "-Command",
            f"{action}-PnpDevice -FriendlyName '*{name}*' -Confirm:$false"],
            capture_output=True, text=True, encoding="utf-8", errors="replace")
        print(r.stdout or f"✅ 裝置已{'啟用' if enable else '停用'}：{name}")
    except Exception as e:
        print(f"❌ 裝置操作失敗：{e}")


# ── 網路介面卡控制 ───────────────────────────────────

def netadapter_list():
    try:
        r = subprocess.run(["powershell.exe", "-Command",
            "Get-NetAdapter | Select-Object Name,Status,MacAddress,LinkSpeed | Format-Table -AutoSize"],
            capture_output=True, text=True, encoding="utf-8", errors="replace")
        print(r.stdout.strip())
    except Exception as e:
        print(f"❌ 查詢失敗：{e}")

def netadapter_toggle(name: str, enable: bool = True):
    try:
        action = "Enable" if enable else "Disable"
        r = subprocess.run(["powershell.exe", "-Command", f"{action}-NetAdapter -Name '{name}' -Confirm:$false"],
            capture_output=True, text=True, encoding="utf-8", errors="replace")
        print(r.stdout or f"✅ 網路卡已{'啟用' if enable else '停用'}：{name}")
    except Exception as e:
        print(f"❌ 網路卡操作失敗：{e}")


# ── DNS/IP 設定 ──────────────────────────────────────

def dns_get():
    try:
        r = subprocess.run(["powershell.exe", "-Command",
            "Get-DnsClientServerAddress | Select-Object InterfaceAlias,ServerAddresses | Format-Table -AutoSize"],
            capture_output=True, text=True, encoding="utf-8", errors="replace")
        print(r.stdout.strip())
    except Exception as e:
        print(f"❌ 查詢失敗：{e}")

def dns_set(interface: str, dns1: str, dns2: str = ""):
    try:
        servers = f"'{dns1}'" + (f",'{dns2}'" if dns2 else "")
        r = subprocess.run(["powershell.exe", "-Command",
            f"Set-DnsClientServerAddress -InterfaceAlias '{interface}' -ServerAddresses ({servers})"],
            capture_output=True, text=True, encoding="utf-8", errors="replace")
        print(r.stdout or f"✅ DNS 已設定：{dns1}{', '+dns2 if dns2 else ''}")
    except Exception as e:
        print(f"❌ DNS 設定失敗：{e}")

def ip_config(interface: str, ip: str, mask: str = "255.255.255.0", gateway: str = ""):
    try:
        ps = f"New-NetIPAddress -InterfaceAlias '{interface}' -IPAddress '{ip}' -PrefixLength 24 -DefaultGateway '{gateway}'"
        r = subprocess.run(["powershell.exe", "-Command", ps],
            capture_output=True, text=True, encoding="utf-8", errors="replace")
        print(r.stdout or f"✅ IP 已設定：{ip}")
    except Exception as e:
        print(f"❌ IP 設定失敗：{e}")


# ── Hosts 檔案編輯 ───────────────────────────────────

HOSTS_PATH = r"C:\Windows\System32\drivers\etc\hosts"

def hosts_list():
    try:
        content = Path(HOSTS_PATH).read_text(encoding="utf-8", errors="replace")
        lines = [l for l in content.splitlines() if l.strip() and not l.strip().startswith("#")]
        print("\n".join(lines) or "（無自訂 hosts）")
    except Exception as e:
        print(f"❌ 讀取 hosts 失敗：{e}")

def hosts_add(ip: str, domain: str):
    try:
        content = Path(HOSTS_PATH).read_text(encoding="utf-8", errors="replace")
        entry = f"\n{ip}\t{domain}"
        if domain in content:
            print(f"⚠️ {domain} 已存在 hosts 中")
            return
        with open(HOSTS_PATH, "a", encoding="utf-8") as f:
            f.write(entry)
        print(f"✅ 已新增：{ip} → {domain}")
    except Exception as e:
        print(f"❌ 新增 hosts 失敗（需管理員權限）：{e}")

def hosts_remove(domain: str):
    try:
        content = Path(HOSTS_PATH).read_text(encoding="utf-8", errors="replace")
        lines = [l for l in content.splitlines() if domain not in l]
        Path(HOSTS_PATH).write_text("\n".join(lines), encoding="utf-8")
        print(f"✅ 已移除：{domain}")
    except Exception as e:
        print(f"❌ 移除 hosts 失敗（需管理員權限）：{e}")


# ── 網路流量監控 ─────────────────────────────────────

def net_traffic(duration: int = 5):
    try:
        import psutil, time as t
        before = psutil.net_io_counters(pernic=False)
        t.sleep(duration)
        after = psutil.net_io_counters(pernic=False)
        sent = (after.bytes_sent - before.bytes_sent) / duration / 1024
        recv = (after.bytes_recv - before.bytes_recv) / duration / 1024
        print(f"📡 網路流量（{duration}秒平均）\n上傳：{sent:.1f} KB/s\n下載：{recv:.1f} KB/s")
        print("\n各網路介面：")
        per_nic = psutil.net_io_counters(pernic=True)
        for name, stats in per_nic.items():
            print(f"  {name}: ↑{stats.bytes_sent/1024/1024:.1f}MB ↓{stats.bytes_recv/1024/1024:.1f}MB")
    except Exception as e:
        print(f"❌ 網路流量監控失敗：{e}")


# ── 條件式自動化 ─────────────────────────────────────

def if_then(condition_type: str, condition_value: str, action_cmd: str, duration: float = 300.0):
    """
    condition_type: cpu_above / mem_above / file_exists / time_is / process_running
    condition_value: 數值或字串
    action_cmd: shell 指令
    """
    try:
        import psutil, time as t
        print(f"🤖 條件式自動化啟動：if [{condition_type}={condition_value}] → [{action_cmd}]")
        end = t.time() + duration
        triggered = False
        while t.time() < end:
            met = False
            if condition_type == "cpu_above":
                met = psutil.cpu_percent(interval=1) > float(condition_value)
            elif condition_type == "mem_above":
                met = psutil.virtual_memory().percent > float(condition_value)
            elif condition_type == "file_exists":
                met = Path(condition_value).exists()
            elif condition_type == "time_is":
                met = dt.datetime.now().strftime("%H:%M") == condition_value
            elif condition_type == "process_running":
                met = any(p.name().lower() == condition_value.lower() for p in psutil.process_iter())
            if met and not triggered:
                subprocess.run(action_cmd, shell=True)
                triggered = True
                print(f"✅ [{dt.datetime.now().strftime('%H:%M:%S')}] 條件觸發，已執行：{action_cmd}")
                break
            t.sleep(2)
        if not triggered:
            print(f"⏰ 監控結束，條件未觸發")
    except Exception as e:
        print(f"❌ 條件式自動化失敗：{e}")


# ── 多視窗排列 ───────────────────────────────────────

def window_arrange(layout: str = "side_by_side"):
    """layout: side_by_side / quad / stack / maximize_all"""
    try:
        import win32gui, win32con
        sw = pyautogui.size().width
        sh = pyautogui.size().height
        hwnds = []
        def _enum(hwnd, _):
            if win32gui.IsWindowVisible(hwnd) and win32gui.GetWindowText(hwnd):
                hwnds.append(hwnd)
        win32gui.EnumWindows(_enum, None)
        hwnds = [h for h in hwnds if win32gui.GetWindowText(h).strip()][:8]
        if layout == "side_by_side" and len(hwnds) >= 2:
            w = sw // 2
            win32gui.MoveWindow(hwnds[0], 0, 0, w, sh, True)
            win32gui.MoveWindow(hwnds[1], w, 0, w, sh, True)
            print(f"✅ 左右分割：{win32gui.GetWindowText(hwnds[0])} | {win32gui.GetWindowText(hwnds[1])}")
        elif layout == "quad" and len(hwnds) >= 4:
            w, h = sw // 2, sh // 2
            positions = [(0,0),(w,0),(0,h),(w,h)]
            for i, (x,y) in enumerate(positions):
                win32gui.MoveWindow(hwnds[i], x, y, w, h, True)
            print(f"✅ 四格排列完成")
        elif layout == "stack":
            h = sh // len(hwnds)
            for i, hwnd in enumerate(hwnds[:4]):
                win32gui.MoveWindow(hwnd, 0, i * h, sw, h, True)
            print(f"✅ 堆疊排列完成")
        elif layout == "maximize_all":
            for hwnd in hwnds:
                win32gui.ShowWindow(hwnd, win32con.SW_MAXIMIZE)
            print(f"✅ 已最大化 {len(hwnds)} 個視窗")
    except Exception as e:
        print(f"❌ 視窗排列失敗：{e}")


# ── 指定區域 OCR ─────────────────────────────────────

def region_ocr(x: int, y: int, w: int, h: int, lang: str = "ch_tra"):
    try:
        import easyocr, numpy as np
        from PIL import Image
        img = pyautogui.screenshot(region=(x, y, w, h))
        reader = easyocr.Reader([lang, "en"], gpu=False)
        results = reader.readtext(np.array(img))
        text = "\n".join(r[1] for r in results)
        print(text or "（未辨識到文字）")
    except Exception as e:
        print(f"❌ 區域 OCR 失敗：{e}")


# ── 指定視窗截圖 ─────────────────────────────────────

def window_screenshot(title_keyword: str, output: str = ""):
    try:
        import win32gui, win32ui, win32con
        from ctypes import windll
        import numpy as np
        from PIL import Image
        hwnd = None
        def _find(h, _):
            nonlocal hwnd
            if title_keyword.lower() in win32gui.GetWindowText(h).lower():
                hwnd = h
        win32gui.EnumWindows(_find, None)
        if not hwnd:
            print(f"❌ 找不到視窗：{title_keyword}")
            return
        win32gui.SetForegroundWindow(hwnd)
        rect = win32gui.GetWindowRect(hwnd)
        x, y, x2, y2 = rect
        w, h = x2 - x, y2 - y
        hwndDC = win32gui.GetWindowDC(hwnd)
        mfcDC = win32ui.CreateDCFromHandle(hwndDC)
        saveDC = mfcDC.CreateCompatibleDC()
        saveBitMap = win32ui.CreateBitmap()
        saveBitMap.CreateCompatibleBitmap(mfcDC, w, h)
        saveDC.SelectObject(saveBitMap)
        windll.user32.PrintWindow(hwnd, saveDC.GetSafeHdc(), 3)
        bmpinfo = saveBitMap.GetInfo()
        bmpstr = saveBitMap.GetBitmapBits(True)
        img = Image.frombuffer("RGB", (bmpinfo["bmWidth"], bmpinfo["bmHeight"]), bmpstr, "raw", "BGRX", 0, 1)
        out = output or str(Path.home() / "Desktop" / f"win_{title_keyword[:10]}_{dt.datetime.now().strftime('%H%M%S')}.png")
        img.save(out)
        win32gui.DeleteObject(saveBitMap.GetHandle())
        saveDC.DeleteDC(); mfcDC.DeleteDC()
        win32gui.ReleaseDC(hwnd, hwndDC)
        print(f"✅ 視窗截圖已存：{out}")
    except Exception as e:
        print(f"❌ 視窗截圖失敗：{e}")


# ── Wave 13：檔案監聽/像素監控/AI物件偵測/滑鼠錄製/ADB/WiFi熱點/OneDrive/FTP/WSL/Hyper-V/檔案diff/螢幕串流 ──

def file_watcher(path: str, events: str = "all", command: str = "", duration: float = 3600.0):
    from watchdog.observers import Observer
    from watchdog.events import FileSystemEventHandler
    import time
    class _H(FileSystemEventHandler):
        def _h(self, e, t):
            if events != "all" and t not in events: return
            print(f"📁 {t}：{e.src_path}")
            if command: import subprocess; subprocess.Popen(command.replace("{path}", e.src_path), shell=True)
        def on_created(self, e): self._h(e, "created")
        def on_modified(self, e): self._h(e, "modified")
        def on_deleted(self, e): self._h(e, "deleted")
    obs = Observer(); obs.schedule(_H(), path, recursive=True); obs.start()
    print(f"👁️ 監聽 {path} 中（{duration}s）...")
    try: time.sleep(duration)
    except KeyboardInterrupt: pass
    finally: obs.stop(); obs.join()

def pixel_watch(x: int, y: int, tolerance: int = 10, interval: float = 1.0, duration: float = 60.0, command: str = ""):
    import pyautogui, time, subprocess
    sc = pyautogui.screenshot(); r0,g0,b0 = sc.getpixel((x,y))[:3]
    print(f"🎨 開始監控({x},{y}) 初始顏色：#{r0:02X}{g0:02X}{b0:02X}")
    end = time.time() + duration
    while time.time() < end:
        sc = pyautogui.screenshot(); r,g,b = sc.getpixel((x,y))[:3]
        if abs(r-r0)+abs(g-g0)+abs(b-b0) > tolerance*3:
            print(f"🔔 顏色變化！#{r:02X}{g:02X}{b:02X}")
            if command: subprocess.Popen(command, shell=True)
            r0,g0,b0 = r,g,b
        time.sleep(interval)

def object_detect(target: str, action: str = "find", region: str = ""):
    import pyautogui, anthropic, base64, io as _io, json, re, os
    reg = None
    if region:
        parts = [int(v) for v in region.split(",")]
        if len(parts) == 4: reg = tuple(parts)
    sc = pyautogui.screenshot(region=reg)
    buf = _io.BytesIO(); sc.save(buf, format="PNG")
    img_b64 = base64.standard_b64encode(buf.getvalue()).decode()
    ox = reg[0] if reg else 0; oy = reg[1] if reg else 0
    cl = anthropic.Anthropic(api_key=os.getenv("ANTHROPIC_API_KEY"))
    resp = cl.messages.create(model="claude-sonnet-4-6", max_tokens=256, messages=[{"role":"user","content":[
        {"type":"image","source":{"type":"base64","media_type":"image/png","data":img_b64}},
        {"type":"text","text":f"找「{target}」位置，回 JSON：{{\"found\":bool,\"x\":int,\"y\":int,\"description\":\"\"}}"}
    ]}])
    m = re.search(r'\{.*\}', resp.content[0].text, re.DOTALL)
    if not m: print("⚠️ 無法解析"); return
    res = json.loads(m.group())
    if not res.get("found"): print(f"⚠️ 未找到：{target}"); return
    ax, ay = res["x"]+ox, res["y"]+oy
    if action == "click": pyautogui.click(ax, ay); print(f"✅ 已點擊「{target}」({ax},{ay})")
    elif action == "double_click": pyautogui.doubleClick(ax, ay); print(f"✅ 已雙擊({ax},{ay})")
    else: print(f"✅ 找到「{target}」：({ax},{ay}) {res.get('description','')}")

def mouse_record(action: str, name: str = "", duration: float = 10.0, repeat: int = 1, speed: float = 1.0):
    import json, time
    f = Path.home() / ".claude_mouse_macros.json"
    store = json.loads(f.read_text()) if f.exists() else {}
    if action == "list":
        print("\n".join(f"- {k}（{len(v)}事件）" for k,v in store.items()) if store else "⚠️ 無已儲存巨集")
    elif action == "delete":
        if name in store: del store[name]; f.write_text(json.dumps(store)); print(f"✅ 已刪除：{name}")
        else: print(f"⚠️ 找不到：{name}")
    elif action == "start":
        from pynput import mouse as m
        events = []; t0 = time.time()
        def on_move(x,y): events.append({"t":time.time()-t0,"type":"move","x":x,"y":y})
        def on_click(x,y,btn,pressed): events.append({"t":time.time()-t0,"type":"click","x":x,"y":y,"btn":str(btn),"pressed":pressed})
        def on_scroll(x,y,dx,dy): events.append({"t":time.time()-t0,"type":"scroll","x":x,"y":y,"dx":dx,"dy":dy})
        lis = m.Listener(on_move=on_move,on_click=on_click,on_scroll=on_scroll)
        lis.start(); time.sleep(duration); lis.stop()
        store[name] = events; f.write_text(json.dumps(store))
        print(f"✅ 已錄製 '{name}'（{len(events)}事件）")
    elif action == "play":
        if name not in store: print(f"⚠️ 找不到：{name}"); return
        from pynput import mouse as m
        ctrl = m.Controller(); evts = store[name]
        for _ in range(repeat):
            prev = 0.0
            for e in evts:
                d = (e["t"]-prev)/speed; time.sleep(min(max(d,0),0.5)); prev=e["t"]
                if e["type"]=="move": ctrl.position=(e["x"],e["y"])
                elif e["type"]=="click":
                    btn=m.Button.left if "left" in e["btn"] else m.Button.right
                    ctrl.press(btn) if e["pressed"] else ctrl.release(btn)
                elif e["type"]=="scroll": ctrl.scroll(e["dx"],e["dy"])
        print(f"✅ '{name}' 回放 {repeat} 次完成")

def adb(action: str, x: int=0, y: int=0, x2: int=0, y2: int=0, text: str="", path: str="", remote: str="", package: str="", command: str="", device: str=""):
    import subprocess
    p = ["adb"] + (["-s",device] if device else [])
    try:
        if action == "devices":
            r = subprocess.run(["adb","devices","-l"],capture_output=True,text=True); print(r.stdout.strip())
        elif action == "screenshot":
            out = path or str(Path.home()/"Desktop"/f"adb_{dt.datetime.now().strftime('%H%M%S')}.png")
            subprocess.run(p+["shell","screencap","-p","/sdcard/screen.png"],capture_output=True)
            subprocess.run(p+["pull","/sdcard/screen.png",out],capture_output=True); print(f"✅ 截圖：{out}")
        elif action == "tap":
            subprocess.run(p+["shell","input","tap",str(x),str(y)],capture_output=True); print(f"✅ 點擊({x},{y})")
        elif action == "swipe":
            subprocess.run(p+["shell","input","swipe",str(x),str(y),str(x2),str(y2),"300"],capture_output=True); print(f"✅ 滑動")
        elif action == "type":
            subprocess.run(p+["shell","input","text",text.replace(" ","%s")],capture_output=True); print(f"✅ 輸入：{text}")
        elif action == "key":
            subprocess.run(p+["shell","input","keyevent",text],capture_output=True); print(f"✅ 按鍵：{text}")
        elif action == "install":
            r=subprocess.run(p+["install","-r",path],capture_output=True,text=True); print(r.stdout.strip())
        elif action == "push":
            r=subprocess.run(p+["push",path,remote or "/sdcard/"],capture_output=True,text=True); print(r.stdout.strip())
        elif action == "pull":
            out=path or str(Path.home()/"Desktop"/Path(remote).name)
            r=subprocess.run(p+["pull",remote,out],capture_output=True,text=True); print(r.stdout.strip())
        elif action == "shell":
            r=subprocess.run(p+["shell",command],capture_output=True,text=True); print(r.stdout.strip())
        elif action == "start_app":
            subprocess.run(p+["shell","monkey","-p",package,"-c","android.intent.category.LAUNCHER","1"],capture_output=True); print(f"✅ 啟動：{package}")
        elif action == "stop_app":
            subprocess.run(p+["shell","am","force-stop",package],capture_output=True); print(f"✅ 停止：{package}")
    except Exception as e: print(f"❌ ADB失敗：{e}")

def wifi_hotspot(action: str, ssid: str="", password: str=""):
    import subprocess
    try:
        if action == "set":
            r=subprocess.run(["netsh","wlan","set","hostednetwork",f"mode=allow",f"ssid={ssid}",f"key={password}"],capture_output=True,text=True,encoding="utf-8",errors="replace"); print(r.stdout.strip())
        elif action == "start":
            r=subprocess.run(["netsh","wlan","start","hostednetwork"],capture_output=True,text=True,encoding="utf-8",errors="replace"); print(r.stdout.strip())
        elif action == "stop":
            r=subprocess.run(["netsh","wlan","stop","hostednetwork"],capture_output=True,text=True,encoding="utf-8",errors="replace"); print(r.stdout.strip())
        elif action == "status":
            r=subprocess.run(["netsh","wlan","show","hostednetwork"],capture_output=True,text=True,encoding="utf-8",errors="replace"); print(r.stdout.strip())
    except Exception as e: print(f"❌ 熱點失敗：{e}")

def onedrive(action: str, path: str="", remote: str=""):
    import os, shutil
    od = os.path.expandvars(r"%USERPROFILE%\OneDrive")
    if not Path(od).exists(): od = str(Path.home()/"OneDrive")
    try:
        if action == "list":
            target = Path(od)/(remote or "")
            for p in sorted(target.iterdir()): print(f"{'📁' if p.is_dir() else '📄'} {p.name}")
        elif action == "upload":
            dest = Path(od)/(remote or Path(path).name); shutil.copy2(path, dest); print(f"✅ 已上傳：{dest}")
        elif action == "download":
            src = Path(od)/remote; out = path or str(Path.home()/"Desktop"/src.name); shutil.copy2(src, out); print(f"✅ 已下載：{out}")
        elif action == "status":
            size = sum(f.stat().st_size for f in Path(od).rglob("*") if f.is_file())/1024/1024/1024; print(f"OneDrive路徑：{od}\n使用：{size:.2f}GB")
        elif action == "open":
            os.startfile(od); print(f"✅ 已開啟：{od}")
    except Exception as e: print(f"❌ OneDrive失敗：{e}")

def ftp(action: str, host: str="", user: str="", password: str="", local: str="", remote: str="", port: int=21):
    from ftplib import FTP
    try:
        f = FTP(); f.connect(host, port, timeout=30); f.login(user, password)
        if action == "list":
            for item in f.nlst(remote or "."): print(item)
        elif action == "upload":
            with open(local,"rb") as fp: f.storbinary(f"STOR {remote or Path(local).name}", fp)
            print(f"✅ 已上傳：{local}")
        elif action == "download":
            out = local or str(Path.home()/"Desktop"/Path(remote).name)
            with open(out,"wb") as fp: f.retrbinary(f"RETR {remote}", fp.write)
            print(f"✅ 已下載：{out}")
        elif action == "delete":
            f.delete(remote); print(f"✅ 已刪除：{remote}")
        elif action == "mkdir":
            f.mkd(remote); print(f"✅ 已建立目錄：{remote}")
        f.quit()
    except Exception as e: print(f"❌ FTP失敗：{e}")

def wsl(action: str, distro: str="", command: str=""):
    import subprocess
    try:
        if action == "list":
            r=subprocess.run(["wsl","--list","--verbose"],capture_output=True,text=True,encoding="utf-16-le",errors="replace"); print(r.stdout.strip())
        elif action == "run":
            cmd=["wsl"]+(["-d",distro] if distro else [])+["--","bash","-c",command]
            r=subprocess.run(cmd,capture_output=True,text=True,encoding="utf-8",errors="replace"); print(r.stdout.strip())
        elif action == "start":
            subprocess.Popen(["wsl"]+(["-d",distro] if distro else [])); print(f"✅ WSL 已啟動")
        elif action == "stop":
            cmd=["wsl","--terminate",distro] if distro else ["wsl","--shutdown"]
            subprocess.run(cmd,capture_output=True); print(f"✅ WSL 已停止")
        elif action == "status":
            r=subprocess.run(["wsl","--status"],capture_output=True,text=True,encoding="utf-16-le",errors="replace"); print(r.stdout.strip())
    except Exception as e: print(f"❌ WSL失敗：{e}")

def hyperv(action: str, name: str="", snapshot: str=""):
    import subprocess
    def ps(cmd):
        r=subprocess.run(["powershell","-Command",cmd],capture_output=True,text=True,encoding="utf-8",errors="replace"); return r.stdout.strip(),r.returncode
    try:
        if action == "list": out,_=ps("Get-VM | Select-Object Name,State | Format-Table -AutoSize"); print(out)
        elif action == "start": ps(f"Start-VM -Name '{name}'"); print(f"✅ 已啟動：{name}")
        elif action == "stop": ps(f"Stop-VM -Name '{name}' -Force"); print(f"✅ 已停止：{name}")
        elif action == "pause": ps(f"Suspend-VM -Name '{name}'"); print(f"✅ 已暫停：{name}")
        elif action == "resume": ps(f"Resume-VM -Name '{name}'"); print(f"✅ 已繼續：{name}")
        elif action == "snapshot":
            sn=snapshot or dt.datetime.now().strftime("snap_%Y%m%d_%H%M%S")
            ps(f"Checkpoint-VM -Name '{name}' -SnapshotName '{sn}'"); print(f"✅ 快照：{sn}")
        elif action == "restore": ps(f"Restore-VMSnapshot -VMName '{name}' -Name '{snapshot}' -Confirm:$false"); print(f"✅ 還原：{snapshot}")
        elif action == "status": out,_=ps(f"Get-VM -Name '{name}' | Format-List"); print(out)
    except Exception as e: print(f"❌ Hyper-V失敗：{e}")

def file_diff(file1: str, file2: str, output: str="", mode: str="unified"):
    import difflib
    try:
        t1=Path(file1).read_text(encoding="utf-8",errors="replace").splitlines(keepends=True)
        t2=Path(file2).read_text(encoding="utf-8",errors="replace").splitlines(keepends=True)
        diff = list(difflib.unified_diff(t1,t2,fromfile=file1,tofile=file2) if mode=="unified" else difflib.context_diff(t1,t2,fromfile=file1,tofile=file2))
        if not diff: print("✅ 兩個檔案完全相同"); return
        result = "".join(diff)
        if output: Path(output).write_text(result,encoding="utf-8"); print(f"✅ diff已存：{output}")
        else: print(result[:3000])
    except Exception as e: print(f"❌ diff失敗：{e}")

def screen_live(fps: float=0.5, duration: float=60.0, quality: int=50):
    import pyautogui, io as _io, time
    interval = 1.0/max(fps,0.1); end=time.time()+duration; count=0
    out_dir=Path.home()/"Desktop"/"screen_live"; out_dir.mkdir(exist_ok=True)
    print(f"📹 開始串流（{fps}FPS，{duration}s）...")
    while time.time()<end:
        sc=pyautogui.screenshot()
        out=str(out_dir/f"frame_{count:04d}.jpg")
        sc.save(out,format="JPEG",quality=quality); count+=1
        time.sleep(interval)
    print(f"✅ 串流完成，{count}張截圖存於：{out_dir}")


# ── Wave 12：AI視覺循環/監控告警/間隔排程/等待文字/瀏覽器進階/語音命令/通知攔截/資料處理/WOL/剪貼簿歷史 ──

def vision_loop(goal: str, max_steps: int = 20, interval: float = 3.0, timeout: float = 120.0):
    import pyautogui, anthropic, base64, io as _io, time, json, re, os
    _client = anthropic.Anthropic(api_key=os.getenv("ANTHROPIC_API_KEY"))
    steps = 0; start = time.time(); log = []
    while steps < max_steps and (time.time() - start) < timeout:
        screenshot = pyautogui.screenshot()
        buf = _io.BytesIO(); screenshot.save(buf, format="PNG")
        img_b64 = base64.standard_b64encode(buf.getvalue()).decode()
        resp = _client.messages.create(model="claude-sonnet-4-6", max_tokens=512, messages=[{"role":"user","content":[
            {"type":"image","source":{"type":"base64","media_type":"image/png","data":img_b64}},
            {"type":"text","text":f"目標：{goal}\n已執行：{log}\n回答 JSON：{{\"done\":bool,\"action\":\"說明\",\"type\":\"click/type/key/wait\",\"x\":0,\"y\":0,\"text\":\"\"}}"}
        ]}])
        m = re.search(r'\{.*\}', resp.content[0].text, re.DOTALL)
        if not m: steps += 1; time.sleep(interval); continue
        act = json.loads(m.group())
        if act.get("done"): print(f"✅ 目標達成！{steps} 步\n" + "\n".join(log)); return
        if act.get("type") == "click" and act.get("x"): pyautogui.click(act["x"], act["y"])
        elif act.get("type") == "type" and act.get("text"): pyautogui.typewrite(act["text"], interval=0.05)
        elif act.get("type") == "key" and act.get("text"): pyautogui.press(act["text"])
        elif act.get("type") == "wait": time.sleep(2)
        log.append(f"步驟{steps+1}: {act.get('action','')}")
        steps += 1; time.sleep(interval)
    print(f"⏳ 完成 {steps} 步\n" + "\n".join(log))

def alert_monitor(condition: str, threshold: str, target: str = "", interval: int = 30, duration: int = 3600):
    import psutil, time
    end = time.time() + duration
    print(f"📊 監控啟動：{condition} {threshold}，每 {interval}s 檢查")
    while time.time() < end:
        try:
            triggered = False; msg = ""
            if condition == "cpu_above":
                v = psutil.cpu_percent(1)
                if v > float(threshold): triggered = True; msg = f"CPU {v:.1f}% > {threshold}%"
            elif condition == "mem_above":
                v = psutil.virtual_memory().percent
                if v > float(threshold): triggered = True; msg = f"記憶體 {v:.1f}% > {threshold}%"
            elif condition == "disk_above":
                v = psutil.disk_usage("/").percent
                if v > float(threshold): triggered = True; msg = f"磁碟 {v:.1f}% > {threshold}%"
            elif condition == "process_missing":
                pnames = [p.name().lower() for p in psutil.process_iter(["name"])]
                if not any(target.lower() in n for n in pnames): triggered = True; msg = f"程序 {target} 已停止"
            elif condition == "screen_text_found":
                import pyautogui, easyocr, tempfile
                reader = easyocr.Reader(["ch_tra","en"], gpu=False)
                screenshot = pyautogui.screenshot()
                tmp = tempfile.mktemp(suffix=".png"); screenshot.save(tmp)
                results = reader.readtext(tmp, detail=0); Path(tmp).unlink(missing_ok=True)
                if target.lower() in " ".join(results).lower(): triggered = True; msg = f"螢幕出現文字：{target}"
            if triggered: print(f"🔔 告警觸發：{msg}")
        except Exception as e: print(f"❌ {e}")
        time.sleep(interval)
    print("✅ 監控結束")

def interval_schedule(command: str, every_minutes: float = 60.0, repeat: int = 0, duration_hours: float = 0.0):
    import subprocess, time
    end = time.time() + duration_hours * 3600 if duration_hours > 0 else float("inf")
    count = 0; max_count = repeat if repeat > 0 else float("inf")
    print(f"⏱️ 間隔排程啟動：{command}，每 {every_minutes} 分鐘")
    while time.time() < end and count < max_count:
        subprocess.Popen(command, shell=True)
        count += 1; print(f"✅ 第 {count} 次執行")
        time.sleep(every_minutes * 60)
    print(f"✅ 排程結束，共執行 {count} 次")

def wait_for_text(text: str, timeout: float = 60.0, interval: float = 2.0, region: str = ""):
    import pyautogui, easyocr, time, tempfile
    reader = easyocr.Reader(["ch_tra","en"], gpu=False)
    start = time.time()
    reg = None
    if region:
        parts = [int(v) for v in region.split(",")]
        if len(parts) == 4: reg = tuple(parts)
    while time.time() - start < timeout:
        screenshot = pyautogui.screenshot(region=reg)
        tmp = tempfile.mktemp(suffix=".png"); screenshot.save(tmp)
        results = reader.readtext(tmp, detail=0); Path(tmp).unlink(missing_ok=True)
        if text.lower() in " ".join(results).lower():
            print(f"✅ 偵測到文字「{text}」（{time.time()-start:.1f}s）"); return
        time.sleep(interval)
    print(f"⏳ 超時，未偵測到「{text}」")

def data_process(action: str, path: str = "", output: str = "", query: str = "", data: str = "", paths: str = ""):
    import json, csv
    try:
        if action == "read_json":
            content = json.loads(Path(path).read_text(encoding="utf-8"))
            if isinstance(content, list): print(f"JSON {len(content)} 筆：\n{json.dumps(content[:5], ensure_ascii=False, indent=2)}")
            else: print(json.dumps(content, ensure_ascii=False, indent=2)[:2000])
        elif action == "write_json":
            Path(output or path).write_text(json.dumps(json.loads(data), ensure_ascii=False, indent=2), encoding="utf-8")
            print(f"✅ JSON 已儲存：{output or path}")
        elif action == "read_csv":
            with open(path, encoding="utf-8-sig", errors="replace") as f:
                rows = list(csv.DictReader(f))
            print(f"CSV {len(rows)} 筆：\n{json.dumps(rows[:5], ensure_ascii=False, indent=2)}")
        elif action == "write_csv":
            obj = json.loads(data)
            with open(output or path, "w", newline="", encoding="utf-8-sig") as f:
                w = csv.DictWriter(f, fieldnames=obj[0].keys()); w.writeheader(); w.writerows(obj)
            print(f"✅ CSV 已儲存：{output or path}")
        elif action == "filter":
            with open(path, encoding="utf-8-sig", errors="replace") as f:
                rows = list(csv.DictReader(f))
            filtered = [r for r in rows if eval(query, {"__builtins__":{}}, r)]
            print(f"過濾 {len(filtered)}/{len(rows)} 筆：\n{json.dumps(filtered[:10], ensure_ascii=False, indent=2)}")
        elif action == "stats":
            with open(path, encoding="utf-8-sig", errors="replace") as f:
                rows = list(csv.DictReader(f))
            for col in rows[0].keys():
                vals = [r[col] for r in rows if r[col]]
                try:
                    nums = [float(v) for v in vals]
                    print(f"{col}: min={min(nums)} max={max(nums)} avg={sum(nums)/len(nums):.2f}")
                except: print(f"{col}: {len(vals)} 筆，{len(set(vals))} 種")
        elif action == "convert":
            ext_in = Path(path).suffix.lower(); ext_out = Path(output).suffix.lower()
            if ext_in == ".json" and ext_out == ".csv":
                obj = json.loads(Path(path).read_text(encoding="utf-8"))
                if not isinstance(obj, list): obj = [obj]
                with open(output, "w", newline="", encoding="utf-8-sig") as f:
                    w = csv.DictWriter(f, fieldnames=obj[0].keys()); w.writeheader(); w.writerows(obj)
            elif ext_in == ".csv" and ext_out == ".json":
                with open(path, encoding="utf-8-sig", errors="replace") as f:
                    rows = list(csv.DictReader(f))
                Path(output).write_text(json.dumps(rows, ensure_ascii=False, indent=2), encoding="utf-8")
            print(f"✅ 已轉換：{output}")
        elif action == "merge":
            all_rows = []
            for p in paths.split(","):
                p = p.strip()
                if p.endswith(".csv"):
                    with open(p, encoding="utf-8-sig", errors="replace") as f: all_rows.extend(list(csv.DictReader(f)))
                elif p.endswith(".json"):
                    obj = json.loads(Path(p).read_text(encoding="utf-8"))
                    if isinstance(obj, list): all_rows.extend(obj)
            out = output or str(Path(paths.split(",")[0].strip()).parent / "merged.csv")
            with open(out, "w", newline="", encoding="utf-8-sig") as f:
                w = csv.DictWriter(f, fieldnames=all_rows[0].keys()); w.writeheader(); w.writerows(all_rows)
            print(f"✅ 合併 {len(all_rows)} 筆 → {out}")
        elif action == "to_table":
            with open(path, encoding="utf-8-sig", errors="replace") as f:
                rows = list(csv.DictReader(f))
            cols = list(rows[0].keys())
            print(" | ".join(cols))
            print("-" * 60)
            for r in rows[:20]: print(" | ".join(str(r.get(c,"")) for c in cols))
    except Exception as e:
        print(f"❌ 資料處理失敗：{e}")

def wake_on_lan(mac: str, broadcast: str = "255.255.255.255", port: int = 9):
    import socket
    try:
        mac_clean = mac.replace(":","").replace("-","")
        magic = b"\xff" * 6 + bytes.fromhex(mac_clean) * 16
        with socket.socket(socket.AF_INET, socket.SOCK_DGRAM) as s:
            s.setsockopt(socket.SOL_SOCKET, socket.SO_BROADCAST, 1)
            s.sendto(magic, (broadcast, port))
        print(f"✅ WOL 封包已送出 → {mac}")
    except Exception as e:
        print(f"❌ WOL 失敗：{e}")

_cb_history = []



# ── Wave 11：防火牆/程序/電源/事件/時間/UI自動化/巨集/顏色/攝影機/多螢幕/印表機/WiFi/代理/鎖定/Defender ──


def process_mgr(action: str, name: str = "", pid: int = None, level: str = "normal"):
    import psutil
    priority_map = {"realtime":psutil.REALTIME_PRIORITY_CLASS,"high":psutil.HIGH_PRIORITY_CLASS,"above_normal":psutil.ABOVE_NORMAL_PRIORITY_CLASS,"normal":psutil.NORMAL_PRIORITY_CLASS,"below_normal":psutil.BELOW_NORMAL_PRIORITY_CLASS,"idle":psutil.IDLE_PRIORITY_CLASS}
    try:
        if action == "list":
            procs = sorted(psutil.process_iter(["pid","name","cpu_percent","memory_info"]), key=lambda p: p.info["cpu_percent"] or 0, reverse=True)
            print(f"{'PID':>6} {'CPU%':>6} {'MEM MB':>8}  名稱")
            for p in procs[:25]:
                mem = (p.info["memory_info"].rss//1024//1024) if p.info["memory_info"] else 0
                print(f"{p.info['pid']:>6} {p.info['cpu_percent'] or 0:>6.1f} {mem:>8}  {p.info['name']}")
        elif action == "search":
            found = [p for p in psutil.process_iter(["pid","name","cpu_percent","memory_info"]) if name.lower() in p.info["name"].lower()]
            if not found: print(f"⚠️ 找不到：{name}"); return
            for p in found:
                print(f"PID:{p.info['pid']} CPU:{p.info['cpu_percent']}% MEM:{p.info['memory_info'].rss//1024//1024}MB {p.info['name']}")
        elif action == "kill":
            targets = [psutil.Process(int(pid))] if pid else [p for p in psutil.process_iter(["pid","name"]) if name.lower() in p.info["name"].lower()]
            if not targets: print(f"⚠️ 找不到：{name}"); return
            for p in targets: p.kill()
            print(f"✅ 已終止 {len(targets)} 個程序：{name or pid}")
        elif action == "priority":
            p = psutil.Process(int(pid)) if pid else next((x for x in psutil.process_iter(["pid","name"]) if name.lower() in x.info["name"].lower()), None)
            if not p: print(f"⚠️ 找不到：{name}"); return
            p.nice(priority_map.get(level, psutil.NORMAL_PRIORITY_CLASS))
            print(f"✅ 已設定 PID {p.pid} 優先權為 {level}")
    except Exception as e:
        print(f"❌ 程序管理失敗：{e}")

def power_plan(action: str, plan: str = "balanced"):
    import subprocess
    plan_guids = {"balanced":"381b4222-f694-41f0-9685-ff5bb260df2e","high_performance":"8c5e7fda-e8bf-4a96-9a85-a6e23a8c635c","power_saver":"a1841308-3541-4fab-bc81-f71556f20b4a"}
    try:
        if action == "list":
            r = subprocess.run(["powercfg","/list"], capture_output=True, text=True, encoding="utf-8", errors="replace")
            print(r.stdout.strip())
        elif action == "get":
            r = subprocess.run(["powercfg","/getactivescheme"], capture_output=True, text=True, encoding="utf-8", errors="replace")
            print(r.stdout.strip())
        elif action == "set":
            guid = plan_guids.get(plan)
            if not guid: print(f"⚠️ 未知計畫：{plan}"); return
            subprocess.run(["powercfg","/setactive",guid], capture_output=True)
            print(f"✅ 電源計畫：{plan}")
    except Exception as e:
        print(f"❌ 電源計畫失敗：{e}")


def datetime_config(action: str, timezone: str = "", datetime_str: str = ""):
    import subprocess
    try:
        if action == "get":
            r = subprocess.run(["powershell","-Command","Get-Date | Format-List; (Get-TimeZone).DisplayName"], capture_output=True, text=True, encoding="utf-8", errors="replace")
            print(r.stdout.strip())
        elif action == "sync":
            subprocess.run(["powershell","-Command","Start-Service w32tm -ErrorAction SilentlyContinue; w32tm /resync /force"], capture_output=True)
            print("✅ 已同步網路時間")
        elif action == "set_timezone":
            r = subprocess.run(["powershell","-Command",f"Set-TimeZone -Id '{timezone}'"], capture_output=True, text=True)
            print(f"✅ 時區：{timezone}" if r.returncode==0 else f"❌ 失敗：{r.stderr.strip()}")
        elif action == "set_time":
            r = subprocess.run(["powershell","-Command",f"Set-Date -Date '{datetime_str}'"], capture_output=True, text=True)
            print(f"✅ 時間：{datetime_str}" if r.returncode==0 else f"❌ 失敗：{r.stderr.strip()}")
    except Exception as e:
        print(f"❌ 時間設定失敗：{e}")

def ui_auto(action: str, window: str = "", control: str = "", text: str = ""):
    from pywinauto import Desktop
    try:
        desktop = Desktop(backend="uia")
        if action == "get_windows":
            wins = [w.window_text() for w in desktop.windows() if w.window_text()]
            print("\n".join(f"- {w}" for w in wins[:30]))
            return
        win = next((w for w in desktop.windows() if window.lower() in w.window_text().lower()), None)
        if not win: print(f"⚠️ 找不到視窗：{window}"); return
        if action == "find":
            info = [f"[{c.control_type()}] {c.window_text()}" for c in win.descendants() if c.window_text()][:30]
            print("\n".join(info))
        elif action == "read":
            print("\n".join(c.window_text() for c in win.descendants() if c.window_text())[:50])
        elif action == "click":
            c = next((c for c in win.descendants() if control.lower() in c.window_text().lower()), None)
            if c: c.click_input(); print(f"✅ 已點擊：{c.window_text()}")
            else: print(f"⚠️ 找不到：{control}")
        elif action == "type":
            c = next((c for c in win.descendants() if control.lower() in c.window_text().lower() or c.control_type() in ("Edit","Document")), None)
            if c: c.type_keys(text); print(f"✅ 已輸入文字")
            else: print(f"⚠️ 找不到輸入框：{control}")
    except Exception as e:
        print(f"❌ UI 自動化失敗：{e}")

def macro(action: str, name: str = "", repeat: int = 1, duration: float = 10.0):
    import keyboard, json, time
    macro_file = Path.home() / ".claude_macros.json"
    store = json.loads(macro_file.read_text()) if macro_file.exists() else {}
    try:
        if action == "record_start":
            keyboard.start_recording()
            time.sleep(duration)
            recorded = keyboard.stop_recording()
            store[name] = [{"type":e.event_type,"name":e.name,"time":e.time} for e in recorded]
            macro_file.write_text(json.dumps(store))
            print(f"✅ 已錄製巨集 '{name}'（{len(recorded)} 事件）")
        elif action == "play":
            if name not in store: print(f"⚠️ 找不到巨集：{name}"); return
            events = store[name]
            for _ in range(repeat):
                prev = events[0]["time"] if events else 0
                for e in events:
                    time.sleep(min(max(0,e["time"]-prev),0.5)); prev=e["time"]
                    keyboard.press(e["name"]) if e["type"]=="down" else keyboard.release(e["name"])
            print(f"✅ 已回放巨集 '{name}' {repeat} 次")
        elif action == "list":
            print("\n".join(f"- {k}（{len(v)} 事件）" for k,v in store.items()) if store else "⚠️ 無已儲存巨集")
        elif action == "delete":
            if name in store:
                del store[name]; macro_file.write_text(json.dumps(store)); print(f"✅ 已刪除：{name}")
            else: print(f"⚠️ 找不到：{name}")
    except Exception as e:
        print(f"❌ 巨集操作失敗：{e}")

def color_pick(x: int, y: int, action: str = "pick", w: int = 100, h: int = 100):
    import pyautogui
    from collections import Counter
    try:
        screenshot = pyautogui.screenshot()
        if action == "pick":
            r,g,b = screenshot.getpixel((x,y))[:3]
            print(f"🎨 ({x},{y}) RGB:({r},{g},{b}) HEX:#{r:02X}{g:02X}{b:02X}")
        elif action == "dominant":
            region = screenshot.crop((x,y,x+w,y+h)).convert("RGB").resize((50,50))
            top5 = Counter(list(region.getdata())).most_common(5)
            for (r,g,b),cnt in top5:
                print(f"RGB({r},{g},{b}) #{r:02X}{g:02X}{b:02X}  出現{cnt}次")
    except Exception as e:
        print(f"❌ 顏色選取失敗：{e}")

def webcam(action: str, device: int = 0, duration: float = 5.0, output: str = ""):
    import cv2
    try:
        if action == "list":
            found = []
            for i in range(5):
                cap = cv2.VideoCapture(i)
                if cap.isOpened(): found.append(f"裝置 {i}"); cap.release()
            print("\n".join(found) if found else "⚠️ 無可用攝影機")
        elif action == "photo":
            cap = cv2.VideoCapture(device)
            if not cap.isOpened(): print(f"❌ 無法開啟攝影機 {device}"); return
            ret, frame = cap.read(); cap.release()
            if not ret: print("❌ 無法拍攝"); return
            out = output or str(Path.home()/"Desktop"/f"webcam_{dt.datetime.now().strftime('%Y%m%d_%H%M%S')}.jpg")
            cv2.imwrite(out, frame); print(f"✅ 已拍照：{out}")
        elif action == "video":
            cap = cv2.VideoCapture(device)
            if not cap.isOpened(): print(f"❌ 無法開啟攝影機 {device}"); return
            out = output or str(Path.home()/"Desktop"/f"webcam_{dt.datetime.now().strftime('%Y%m%d_%H%M%S')}.avi")
            fw,fh = int(cap.get(3)),int(cap.get(4))
            writer = cv2.VideoWriter(out,cv2.VideoWriter_fourcc(*"XVID"),20,(fw,fh))
            import time; end=time.time()+duration
            while time.time()<end:
                ret,frame=cap.read()
                if ret: writer.write(frame)
            cap.release(); writer.release(); print(f"✅ 已錄影：{out}")
    except Exception as e:
        print(f"❌ 攝影機操作失敗：{e}")

def multi_monitor(action: str, monitor: int = 1, window: str = ""):
    import subprocess, win32gui
    try:
        if action == "list":
            r = subprocess.run(["powershell","-Command","Get-CimInstance Win32_DesktopMonitor | Select-Object Name,ScreenWidth,ScreenHeight | Format-Table -AutoSize"], capture_output=True, text=True, encoding="utf-8", errors="replace")
            print(r.stdout.strip())
        elif action in ("extend","clone"):
            subprocess.Popen(["displayswitch.exe", "/extend" if action=="extend" else "/clone"])
            print(f"✅ 螢幕模式：{action}")
        elif action == "move_window":
            import ctypes; sw=ctypes.windll.user32.GetSystemMetrics(0)
            hwnds=[]
            win32gui.EnumWindows(lambda h,l: l.append(h) if win32gui.IsWindowVisible(h) and window.lower() in win32gui.GetWindowText(h).lower() else None, hwnds)
            if not hwnds: print(f"⚠️ 找不到視窗：{window}"); return
            rect=win32gui.GetWindowRect(hwnds[0])
            win32gui.MoveWindow(hwnds[0],sw*(monitor-1)+100,100,rect[2]-rect[0],rect[3]-rect[1],True)
            print(f"✅ 視窗 '{window}' 移至螢幕 {monitor}")
    except Exception as e:
        print(f"❌ 多螢幕管理失敗：{e}")

def printer(action: str, path: str = "", printer_name: str = ""):
    import win32print
    try:
        if action == "list":
            printers = win32print.EnumPrinters(win32print.PRINTER_ENUM_LOCAL|win32print.PRINTER_ENUM_CONNECTIONS)
            default = win32print.GetDefaultPrinter()
            print(f"預設：{default}")
            for p in printers: print(f"- {p[2]}")
        elif action == "print":
            import win32api
            pname = printer_name or win32print.GetDefaultPrinter()
            win32api.ShellExecute(0,"print",path,f'/d:"{pname}"',".",0)
            print(f"✅ 已傳送列印：{path}")
        elif action == "queue":
            pname = printer_name or win32print.GetDefaultPrinter()
            h = win32print.OpenPrinter(pname)
            jobs = win32print.EnumJobs(h,0,-1,1)
            win32print.ClosePrinter(h)
            print("佇列為空" if not jobs else "\n".join(f"工作{j['JobId']}：{j['pDocument']}" for j in jobs))
        elif action == "clear_queue":
            pname = printer_name or win32print.GetDefaultPrinter()
            h = win32print.OpenPrinter(pname)
            for j in win32print.EnumJobs(h,0,-1,1):
                win32print.SetJob(h,j["JobId"],0,None,win32print.JOB_CONTROL_DELETE)
            win32print.ClosePrinter(h); print(f"✅ 佇列已清除")
        elif action == "set_default":
            win32print.SetDefaultPrinter(printer_name); print(f"✅ 預設印表機：{printer_name}")
    except Exception as e:
        print(f"❌ 印表機操作失敗：{e}")

def wifi(action: str, ssid: str = "", password: str = ""):
    import subprocess
    try:
        cmd_map = {"scan":["netsh","wlan","show","networks","mode=Bssid"],"status":["netsh","wlan","show","interfaces"],"saved":["netsh","wlan","show","profiles"],"disconnect":["netsh","wlan","disconnect"]}
        if action in cmd_map:
            r = subprocess.run(cmd_map[action], capture_output=True, text=True, encoding="utf-8", errors="replace")
            print(r.stdout.strip()[:2000])
        elif action == "password":
            r = subprocess.run(["netsh","wlan","show","profile",f"name={ssid}","key=clear"], capture_output=True, text=True, encoding="utf-8", errors="replace")
            print(r.stdout.strip())
        elif action == "connect":
            r = subprocess.run(["netsh","wlan","connect",f"name={ssid}"], capture_output=True, text=True, encoding="utf-8", errors="replace")
            print(r.stdout.strip())
    except Exception as e:
        print(f"❌ WiFi 操作失敗：{e}")

def proxy(action: str, host: str = ""):
    import winreg
    key_path = r"Software\Microsoft\Windows\CurrentVersion\Internet Settings"
    try:
        with winreg.OpenKey(winreg.HKEY_CURRENT_USER, key_path, 0, winreg.KEY_READ|winreg.KEY_SET_VALUE) as k:
            if action == "get":
                try:
                    enabled = winreg.QueryValueEx(k,"ProxyEnable")[0]
                    server = winreg.QueryValueEx(k,"ProxyServer")[0]
                    print(f"代理：{'啟用' if enabled else '停用'}  伺服器：{server}")
                except: print("代理：未設定")
            elif action == "set":
                winreg.SetValueEx(k,"ProxyEnable",0,winreg.REG_DWORD,1)
                winreg.SetValueEx(k,"ProxyServer",0,winreg.REG_SZ,host)
                print(f"✅ 代理已設定：{host}")
            elif action == "disable":
                winreg.SetValueEx(k,"ProxyEnable",0,winreg.REG_DWORD,0)
                print("✅ 代理已停用")
    except Exception as e:
        print(f"❌ 代理設定失敗：{e}")

def lock_screen(action: str = "lock"):
    import subprocess
    try:
        if action == "lock":
            subprocess.Popen(["rundll32.exe","user32.dll,LockWorkStation"]); print("🔒 螢幕已鎖定")
        elif action == "logoff":
            subprocess.run(["shutdown","/l"],capture_output=True); print("✅ 已登出")
        elif action == "switch_user":
            subprocess.Popen(["tsdiscon.exe"]); print("✅ 已切換使用者")
    except Exception as e:
        print(f"❌ 鎖定/登出失敗：{e}")

def defender(action: str, path: str = ""):
    import subprocess
    try:
        if action == "status":
            r = subprocess.run(["powershell","-Command","Get-MpComputerStatus | Select-Object AMRunningMode,RealTimeProtectionEnabled,AntivirusSignatureLastUpdated | Format-List"], capture_output=True, text=True, encoding="utf-8", errors="replace")
            print(r.stdout.strip())
        elif action in ("quick_scan","full_scan"):
            scan_type = "QuickScan" if action=="quick_scan" else "FullScan"
            subprocess.Popen(["powershell","-Command",f"Start-MpScan -ScanType {scan_type}"])
            print(f"🛡️ {action} 已啟動（背景執行中）")
        elif action == "threats":
            r = subprocess.run(["powershell","-Command","Get-MpThreatDetection | Select-Object ThreatID,Resources,ActionSuccess | Format-Table -AutoSize"], capture_output=True, text=True, encoding="utf-8", errors="replace")
            print(r.stdout.strip() or "✅ 無威脅記錄")
        elif action == "add_exclusion":
            r = subprocess.run(["powershell","-Command",f"Add-MpPreference -ExclusionPath '{path}'"], capture_output=True, text=True)
            print(f"✅ 已新增排除：{path}" if r.returncode==0 else f"❌ 失敗：{r.stderr.strip()}")
        elif action == "remove_exclusion":
            r = subprocess.run(["powershell","-Command",f"Remove-MpPreference -ExclusionPath '{path}'"], capture_output=True, text=True)
            print(f"✅ 已移除排除：{path}" if r.returncode==0 else f"❌ 失敗：{r.stderr.strip()}")
        elif action == "list_exclusions":
            r = subprocess.run(["powershell","-Command","(Get-MpPreference).ExclusionPath"], capture_output=True, text=True, encoding="utf-8", errors="replace")
            print(r.stdout.strip() or "✅ 無排除項目")
    except Exception as e:
        print(f"❌ Defender 操作失敗：{e}")


# ══════════════════════════════════════════════════════════════════════
# EXTRACTED FUNCTIONS FROM bot.py
# These functions were previously imported via 'from bot import ...'
# Now defined locally for standalone operation.
# ══════════════════════════════════════════════════════════════════════

# ── Global variables needed by extracted functions ──
BOT_DIR = Path("C:/Users/blue_/claude-telegram-bot")
DB_PATH = BOT_DIR / "bot_data.db"
MSG_LOG = BOT_DIR / "msg.log"
_db_conn = None
_db_lock = threading.Lock()
_xtts_server_started = False

def _ensure_xtts_server():
    try:
        r = requests.get("http://localhost:8020/docs", timeout=2)
        return r.status_code == 200
    except Exception:
        return False

def _xtts_generate_wav(text):
    try:
        r = requests.post("http://localhost:8020/tts_to_audio/",
                          json={"text": text, "speaker_wav": "", "language": "zh-cn"},
                          timeout=30)
        if r.status_code == 200:
            return r.content
    except Exception:
        pass
    return None


# Globals
OWNER_ID = 8362721681
_XTTS_PORT = 5678
_xtts_server_proc = None
# Globals from extracted_desktop_control.py
# NOTE: pyautogui is already imported at line 337, do NOT override with None
_pyautogui_dc = None  # placeholder for desktop control lazy-init if needed
_yolo_model = None


def _action_with_verify(action_fn, monitor: int = 1, max_retries: int = 3, wait: float = 0.8):
    """執行動作 → 比對前後截圖 → 沒變化就重試
    action_fn: 無參數的 callable，執行實際動作（如點擊）
    回傳 True 如果畫面有變化，False 如果重試後仍無變化
    """
    import time, numpy as np
    for attempt in range(max_retries):
        img_before, _, _ = _cap_monitor_logical(monitor)
        arr_before = np.array(img_before)
        action_fn()
        time.sleep(wait)
        img_after, _, _ = _cap_monitor_logical(monitor)
        arr_after = np.array(img_after)
        if arr_before.shape == arr_after.shape:
            diff = np.mean(np.abs(arr_before.astype(float) - arr_after.astype(float)))
            diff_pct = diff / 255 * 100
            if diff_pct > 0.5:  # 畫面有變化 = 動作成功
                return True
        # 沒變化，重試
        time.sleep(0.3)
    return False


def _cap_monitor_logical(monitor: int):
    """截圖回傳 (PIL.Image, mon_left_log, mon_top_log) 使用 mss 邏輯座標"""
    import mss
    from PIL import Image as _PI
    with mss.mss() as s:
        mon = s.monitors[monitor] if monitor < len(s.monitors) else s.monitors[1]
        shot = s.grab(mon)
        img = _PI.frombytes("RGB", shot.size, shot.bgra, "raw", "BGRX")
    return img, mon["left"], mon["top"]




def _get_db() -> sqlite3.Connection:
    global _db_conn
    if _db_conn is None:
        _db_conn = sqlite3.connect(DB_PATH, check_same_thread=False)
        _db_conn.execute("PRAGMA journal_mode=WAL")
        _db_conn.execute("PRAGMA synchronous=NORMAL")
        _db_conn.execute("PRAGMA cache_size=2000")
    return _db_conn


def _get_monitors():
    """取得所有螢幕資訊，回傳 list of dict（index 1-based，跳過 monitor[0] 全域合併區）"""
    import mss
    with mss.mss() as sct:
        return sct.monitors[1:]  # index 0 是全域，1~ 才是實體螢幕


def _get_virtual_desktop():
    """回傳 (vl, vt, vw, vh, is_phys, sx, sy)
    vl/vt/vw/vh = 物理虛擬桌面，sx/sy = 邏輯→物理縮放比"""
    import ctypes as _c, mss as _m
    u32 = _c.windll.user32
    old = u32.SetThreadDpiAwarenessContext(_c.c_void_p(-1))
    vl = u32.GetSystemMetrics(76); vt = u32.GetSystemMetrics(77)
    vw = u32.GetSystemMetrics(78); vh = u32.GetSystemMetrics(79)
    u32.SetThreadDpiAwarenessContext(old)
    with _m.mss() as s:
        ms = s.monitors[1:]
        vl_log = min(m["left"] for m in ms)
        vt_log = min(m["top"] for m in ms)
        vw_log = max(m["left"]+m["width"] for m in ms) - vl_log
        vh_log = max(m["top"]+m["height"] for m in ms) - vt_log
    is_phys = vw > 4000
    sx = vw / vw_log if is_phys and vw_log > 0 else 1.0
    sy = vh / vh_log if is_phys and vh_log > 0 else 1.0
    return vl, vt, vw, vh, is_phys, sx, sy


def _resolve_coords(x, y, monitor_idx):
    """將相對於指定螢幕的座標轉換成全域座標"""
    if monitor_idx is None:
        return x, y
    monitors = _get_monitors()
    if 1 <= monitor_idx <= len(monitors):
        m = monitors[monitor_idx - 1]
        return m["left"] + x, m["top"] + y
    return x, y


def _save_training_sample(img, description: str, rx: int, ry: int):
    """每次 vision_locate 成功時自動保存截圖+標記，累積 YOLO 訓練資料"""
    try:
        from pathlib import Path
        import json, time
        data_dir = Path(__file__).parent / "yolo_training_data"
        data_dir.mkdir(exist_ok=True)
        ts = int(time.time() * 1000)
        img.save(data_dir / f"{ts}.jpg", quality=90)
        # YOLO 格式：class x_center y_center width height（歸一化）
        # 暫時用固定 50x50 框，之後可以調整
        w, h = img.width, img.height
        box_w, box_h = 50 / w, 50 / h
        cx, cy = rx / w, ry / h
        with open(data_dir / f"{ts}.txt", "w") as f:
            f.write(f"0 {cx:.6f} {cy:.6f} {box_w:.6f} {box_h:.6f}\n")
        with open(data_dir / f"{ts}.json", "w", encoding="utf-8") as f:
            json.dump({"description": description, "x": rx, "y": ry, "img_w": w, "img_h": h}, f, ensure_ascii=False)
    except Exception:
        pass


def _si_scroll(ax_log: float, ay_log: float, amount: int, direction: str = "down"):
    """SendInput 滾輪，在邏輯座標位置向上/下滾動"""
    import ctypes as _c, ctypes.wintypes as _w, time as _t
    vl, vt, vw, vh, is_phys, sx, sy = _get_virtual_desktop()
    ax = round(ax_log * sx); ay = round(ay_log * sy)
    nx = int((ax - vl) * 65535 // vw)
    ny = int((ay - vt) * 65535 // vh)
    u32 = _c.windll.user32
    class MI(_c.Structure):
        _fields_ = [('dx',_w.LONG),('dy',_w.LONG),('mouseData',_w.DWORD),
                    ('dwFlags',_w.DWORD),('time',_w.DWORD),('dwExtraInfo',_c.POINTER(_c.c_ulong))]
    class U(_c.Union): _fields_ = [('mi', MI)]
    class INP(_c.Structure):
        _anonymous_ = ('u',); _fields_ = [('type',_w.DWORD),('u',U)]
    def _send(flags, dx=0, dy=0, md=0):
        i = INP(0, U(mi=MI(dx, dy, _c.c_ulong(md & 0xFFFFFFFF).value, flags, 0, None)))
        u32.SendInput(1, _c.byref(i), _c.sizeof(i))
    _send(0x0001|0x8000|0x4000, nx, ny)
    _t.sleep(0.15)
    delta = 120 * amount * (1 if direction == "up" else -1)
    for _ in range(max(1, amount // 3)):
        _send(0x0800, md=120 * 3 * (1 if direction == "up" else -1))
        _t.sleep(0.08)


def _si_universal(ax_log: float, ay_log: float, click_type: str = "click"):
    """SendInput 點擊，ax_log/ay_log 為 mss 邏輯座標（支援負值螢幕2）"""
    import ctypes as _c, ctypes.wintypes as _w, time as _t
    vl, vt, vw, vh, is_phys, sx, sy = _get_virtual_desktop()
    ax = round(ax_log * sx); ay = round(ay_log * sy)
    nx = int((ax - vl) * 65535 // vw)
    ny = int((ay - vt) * 65535 // vh)
    u32 = _c.windll.user32
    class MI(_c.Structure):
        _fields_ = [('dx',_w.LONG),('dy',_w.LONG),('mouseData',_w.DWORD),
                    ('dwFlags',_w.DWORD),('time',_w.DWORD),('dwExtraInfo',_c.POINTER(_c.c_ulong))]
    class U(_c.Union): _fields_ = [('mi', MI)]
    class INP(_c.Structure):
        _anonymous_ = ('u',); _fields_ = [('type',_w.DWORD),('u',U)]
    def _send(flags, dx=0, dy=0, md=0):
        i = INP(0, U(mi=MI(dx, dy, md, flags, 0, None)))
        u32.SendInput(1, _c.byref(i), _c.sizeof(i))
    _send(0x0001|0x8000|0x4000, nx, ny)
    _t.sleep(0.25)
    if click_type == "double_click":
        _send(0x0002); _t.sleep(0.05); _send(0x0004); _t.sleep(0.05)
        _send(0x0002); _t.sleep(0.05); _send(0x0004)
    elif click_type == "right_click":
        _send(0x0008); _t.sleep(0.05); _send(0x0010)
    else:
        _send(0x0002); _t.sleep(0.05); _send(0x0004)


def _uia_find_element(description: str, window_title: str = None):
    """用 pywinauto UIA 在前景視窗中搜尋元素，回傳 (center_x, center_y) 或 (None, None)
    比 vision_locate 快 10 倍，不需截圖和 API
    """
    try:
        from pywinauto import Desktop
        import re

        desktop = Desktop(backend="uia")

        # 找目標視窗
        if window_title:
            wins = desktop.windows(title_re=f".*{re.escape(window_title)}.*", visible_only=True)
        else:
            # 用前景視窗
            import win32gui
            hwnd = win32gui.GetForegroundWindow()
            wins = [desktop.window(handle=hwnd)]

        if not wins:
            return None, None

        win = wins[0]
        desc_lower = description.lower()

        # 搜尋策略：先精確匹配，再模糊匹配
        best = None
        best_score = 0

        try:
            # 取得所有子元素（限深度避免太慢）
            children = win.descendants(depth=8)
        except Exception:
            return None, None

        for child in children:
            try:
                name = (child.window_text() or "").strip()
                ctrl_type = (child.friendly_class_name() or "").lower()
                if not name:
                    continue

                name_lower = name.lower()
                score = 0

                # 精確包含
                if desc_lower in name_lower:
                    score = 100
                # 關鍵字匹配
                else:
                    keywords = [w for w in desc_lower.split() if len(w) > 1]
                    if keywords:
                        matched = sum(1 for kw in keywords if kw in name_lower)
                        score = matched / len(keywords) * 80

                # 可點擊元素加分
                if ctrl_type in ("button", "hyperlink", "listitem", "menuitem", "treeitem", "tabitem"):
                    score += 10

                if score > best_score:
                    rect = child.rectangle()
                    cx = (rect.left + rect.right) // 2
                    cy = (rect.top + rect.bottom) // 2
                    # 確保座標在合理範圍
                    if -4000 < cx < 8000 and -2000 < cy < 4000:
                        best = (cx, cy)
                        best_score = score
            except Exception:
                continue

        if best and best_score >= 50:
            return best
        return None, None
    except Exception:
        return None, None


def _vision_find(img, description: str):
    """截圖 → OCR輔助 + Claude Vision → 回傳 (rx, ry) resized圖像像素，找不到回傳 (None, None)"""
    import anthropic, base64, io, json, re
    from PIL import Image as _PI
    ow, oh = img.width, img.height
    # 2048px / quality 92（于晏哥指定的辨識品質標準）
    if img.width > 2048:
        r = 2048 / img.width
        img = img.resize((2048, int(img.height * r)), _PI.LANCZOS)
    scale = ow / img.width
    buf = io.BytesIO(); img.save(buf, format="JPEG", quality=92)
    # 超過 4MB 則降品質重試
    if buf.tell() > 4 * 1024 * 1024:
        buf = io.BytesIO(); img.save(buf, format="JPEG", quality=80)
    b64 = base64.standard_b64encode(buf.getvalue()).decode()

    # OCR 輔助：提取螢幕上的文字供 Claude 參考
    ocr_hint = ""
    try:
        import pytesseract
        ocr_data = pytesseract.image_to_data(img, lang="chi_tra+eng", output_type=pytesseract.Output.DICT)
        ocr_items = []
        for i, txt in enumerate(ocr_data["text"]):
            if txt.strip() and ocr_data["conf"][i] > 40:
                cx = ocr_data["left"][i] + ocr_data["width"][i] // 2
                cy = ocr_data["top"][i] + ocr_data["height"][i] // 2
                ocr_items.append(f"「{txt.strip()}」at({cx},{cy})")
        if ocr_items:
            ocr_hint = "\n\nOCR偵測到的文字及位置：" + "; ".join(ocr_items[:30])
    except Exception:
        pass

    prompt = (
        f"你是專業的螢幕分析師。這是一張電腦螢幕截圖（{img.width}x{img.height}px）。\n"
        f"請找到「{description}」的中心座標。\n"
        f"注意：\n"
        f"- 座標是相對於圖片左上角(0,0)的像素位置\n"
        f"- 仔細區分廣告和真正的內容（廣告通常有「贊助商廣告」「Ad」「探索」字樣）\n"
        f"- 如果有多個匹配，選最相關的那個\n"
        f"- 僅回傳JSON: {{\"x\":整數, \"y\":整數, \"ok\":true/false}}"
        f"{ocr_hint}"
    )

    resp = anthropic.Anthropic().messages.create(
        model="claude-sonnet-4-6", max_tokens=300,
        messages=[{"role": "user", "content": [
            {"type": "image", "source": {"type": "base64", "media_type": "image/jpeg", "data": b64}},
            {"type": "text", "text": prompt}
        ]}]
    )
    m = re.search(r'\{.*?\}', resp.content[0].text, re.DOTALL)
    if not m: return None, None
    d = json.loads(m.group())
    if not d.get("ok", True): return None, None
    rx, ry = int(d["x"] * scale), int(d["y"] * scale)
    # 自動收集訓練資料
    _save_training_sample(img, description, d["x"], d["y"])
    return rx, ry


def _wait_screen_stable(monitor: int = 1, threshold: float = 0.5, timeout: float = 15.0, interval: float = 0.5):
    """連續截圖比對，直到畫面穩定（差異低於 threshold%）或超時"""
    import time, numpy as np
    img_a, _, _ = _cap_monitor_logical(monitor)
    arr_a = np.array(img_a)
    start = time.time()
    while time.time() - start < timeout:
        time.sleep(interval)
        img_b, _, _ = _cap_monitor_logical(monitor)
        arr_b = np.array(img_b)
        if arr_a.shape == arr_b.shape:
            diff = np.mean(np.abs(arr_a.astype(float) - arr_b.astype(float)))
            diff_pct = diff / 255 * 100
            if diff_pct < threshold:
                return True  # 畫面穩定
        arr_a = arr_b
    return False  # 超時




def _yolo_detect(img, conf=0.4):
    """用 YOLO 偵測螢幕上的物件，回傳 [(label, x_center, y_center, w, h, confidence), ...]
    載入 yolo_ui.pt（通用模型或自訂 UI 模型）
    """
    global _yolo_model
    try:
        if _yolo_model is None:
            from pathlib import Path
            model_path = Path(__file__).parent / "yolo_ui.pt"
            if not model_path.exists():
                return []
            from ultralytics import YOLO
            _yolo_model = YOLO(str(model_path))
        import numpy as np
        from PIL import Image as _PI
        if isinstance(img, _PI.Image):
            img_arr = np.array(img)
        else:
            img_arr = img
        results = _yolo_model(img_arr, conf=conf, verbose=False)
        detections = []
        for r in results:
            for box in r.boxes:
                cls_id = int(box.cls[0])
                label = _yolo_model.names[cls_id]
                x1, y1, x2, y2 = box.xyxy[0].tolist()
                cx, cy = int((x1 + x2) / 2), int((y1 + y2) / 2)
                w, h = int(x2 - x1), int(y2 - y1)
                detections.append((label, cx, cy, w, h, float(box.conf[0])))
        return detections
    except Exception as e:
        return []


def analyze_pdf(path: str, max_chars: int = 4000) -> str:
    try:
        import pdfplumber
        with pdfplumber.open(path) as pdf:
            total_pages = len(pdf.pages)
            texts = []
            for page in pdf.pages[:20]:  # 最多讀前 20 頁
                t = page.extract_text()
                if t:
                    texts.append(t)
        full_text = "\n".join(texts)
        if not full_text.strip():
            return "PDF 無法提取文字（可能是掃描圖片 PDF）"
        result = f"📄 PDF 分析（共 {total_pages} 頁）\n\n{full_text[:max_chars]}"
        if len(full_text) > max_chars:
            result += f"\n\n（內容已截斷，共 {len(full_text)} 字）"
        return result
    except Exception as e:
        return f"PDF 讀取失敗：{e}"


def compare_stocks(symbols: list, metrics: list = None) -> str:
    try:
        import yfinance as yf
        if metrics is None or "all" in metrics:
            metrics = ["price", "pe", "roe", "margin", "growth"]

        rows = []
        for sym in symbols[:5]:
            try:
                info = yf.Ticker(sym).info
                hist = yf.Ticker(sym).history(period="1mo")
                ret_1m = ((hist["Close"].iloc[-1] / hist["Close"].iloc[0]) - 1) * 100 if len(hist) > 1 else None
                row = {
                    "symbol": sym,
                    "name": (info.get("shortName") or sym)[:15],
                    "price": info.get("currentPrice") or info.get("regularMarketPrice"),
                    "ret_1m": ret_1m,
                    "pe": info.get("trailingPE"),
                    "pb": info.get("priceToBook"),
                    "roe": info.get("returnOnEquity"),
                    "margin": info.get("profitMargins"),
                    "rev_growth": info.get("revenueGrowth"),
                    "earn_growth": info.get("earningsGrowth"),
                    "div_yield": info.get("dividendYield"),
                    "mkt_cap": info.get("marketCap"),
                }
                rows.append(row)
            except Exception:
                rows.append({"symbol": sym, "name": sym})

        if not rows:
            return "無法取得比較資料"

        lines = [f"📊 股票比較：{' vs '.join(symbols[:5])}\n"]

        def fmt(v, pct=False, mult=100):
            if v is None: return "N/A"
            if pct: return f"{v*mult:+.1f}%"
            return f"{v:.2f}"

        for r in rows:
            sym = r["symbol"]
            name = r.get("name", sym)
            mc = r.get("mkt_cap")
            mc_str = f"{mc/1e12:.2f}T" if mc and mc >= 1e12 else (f"{mc/1e9:.1f}B" if mc else "N/A")
            lines.append(f"── {name} ({sym}) ──")
            if "price" in metrics and r.get("price"):
                lines.append(f"  現價：{r['price']:.2f}　近1月：{fmt(r.get('ret_1m'), False)+'%' if r.get('ret_1m') else 'N/A'}")
            lines.append(f"  市值：{mc_str}")
            if "pe" in metrics:
                lines.append(f"  P/E：{fmt(r.get('pe'))}　P/B：{fmt(r.get('pb'))}")
            if "roe" in metrics and r.get("roe"):
                lines.append(f"  ROE：{r['roe']*100:.1f}%")
            if "margin" in metrics and r.get("margin"):
                lines.append(f"  淨利率：{r['margin']*100:.1f}%")
            if "growth" in metrics:
                lines.append(f"  營收成長：{fmt(r.get('rev_growth'), True)}　獲利成長：{fmt(r.get('earn_growth'), True)}")
            if r.get("div_yield"):
                lines.append(f"  殖利率：{r['div_yield']*100:.2f}%")

        return "\n".join(lines)
    except Exception as e:
        return f"股票比較失敗：{e}"


def execute_adb(action, x=0, y=0, x2=0, y2=0, text="", path="", remote="", package="", command="", device=""):
    try:
        prefix = ["adb"]
        if device: prefix += ["-s", device]
        if action == "devices": r = subprocess.run(["adb", "devices", "-l"], capture_output=True, text=True); return f"📱 ADB 裝置：\n{r.stdout.strip()}"
        elif action == "screenshot":
            out = path or str(Path("C:/Users/blue_/Desktop/測試檔案") / f"adb_{dt.dt.datetime.now().strftime('%H%M%S')}.png")
            subprocess.run(prefix + ["shell", "screencap", "-p", "/sdcard/screen.png"], capture_output=True)
            subprocess.run(prefix + ["pull", "/sdcard/screen.png", out], capture_output=True)
            return f"✅ 手機截圖已存：{out}"
        elif action == "tap": subprocess.run(prefix + ["shell", "input", "tap", str(x), str(y)], capture_output=True); return f"✅ 已點擊手機 ({x},{y})"
        elif action == "swipe": subprocess.run(prefix + ["shell", "input", "swipe", str(x), str(y), str(x2), str(y2), "300"], capture_output=True); return f"✅ 已滑動 ({x},{y})→({x2},{y2})"
        elif action == "type": t = text.replace(" ", "%s"); subprocess.run(prefix + ["shell", "input", "text", t], capture_output=True); return f"✅ 已輸入文字：{text}"
        elif action == "key": subprocess.run(prefix + ["shell", "input", "keyevent", text], capture_output=True); return f"✅ 已按鍵：{text}"
        elif action == "install": r = subprocess.run(prefix + ["install", "-r", path], capture_output=True, text=True); return f"✅ 已安裝：{path}" if r.returncode == 0 else f"❌ 安裝失敗：{r.stderr}"
        elif action == "push": r = subprocess.run(prefix + ["push", path, remote or "/sdcard/"], capture_output=True, text=True); return f"✅ 已上傳" if r.returncode == 0 else f"❌ 失敗：{r.stderr}"
        elif action == "pull":
            out = path or str(Path("C:/Users/blue_/Desktop/測試檔案") / Path(remote).name)
            r = subprocess.run(prefix + ["pull", remote, out], capture_output=True, text=True); return f"✅ 已下載：{remote} → {out}" if r.returncode == 0 else f"❌ 失敗：{r.stderr}"
        elif action == "shell": r = subprocess.run(prefix + ["shell", command], capture_output=True, text=True); return f"📱 ADB Shell：\n{r.stdout.strip()}"
        elif action == "start_app": subprocess.run(prefix + ["shell", "monkey", "-p", package, "-c", "android.intent.category.LAUNCHER", "1"], capture_output=True); return f"✅ 已啟動 App：{package}"
        elif action == "stop_app": subprocess.run(prefix + ["shell", "am", "force-stop", package], capture_output=True); return f"✅ 已停止 App：{package}"
    except Exception as e: return f"❌ ADB 操作失敗：{e}"


def execute_ai_video(prompt: str, provider: str = "replicate",
                     model: str = "", image_url: str = "",
                     duration: float = 5, output: str = "") -> str:
    """用 AI API 生成影片 — provider: replicate / runway / kling"""
    import requests, time, traceback
    from pathlib import Path
    out = output or str(Path("C:/Users/blue_/Desktop/測試檔案") / f"ai_video_{int(time.time())}.mp4")

    def _download(url: str, dest: str) -> str:
        r = requests.get(url, timeout=120, stream=True)
        r.raise_for_status()
        with open(dest, "wb") as f:
            for chunk in r.iter_content(8192):
                f.write(chunk)
        return dest

    try:
        if provider == "replicate":
            api_key = os.getenv("REPLICATE_API_TOKEN", "")
            if not api_key:
                return "❌ 缺少 REPLICATE_API_TOKEN，請在 .env 加入"
            mdl = model or ("stability-ai/stable-video-diffusion" if image_url else "minimax/video-01")
            headers = {"Authorization": f"Bearer {api_key}", "Content-Type": "application/json"}
            inputs = {"prompt": prompt}
            if image_url: inputs["image"] = image_url
            if duration: inputs["duration"] = int(duration)
            r = requests.post(
                f"https://api.replicate.com/v1/models/{mdl}/predictions" if "/" in mdl
                else f"https://api.replicate.com/v1/predictions",
                json={"input": inputs, **({"version": mdl} if "/" not in mdl else {})},
                headers=headers, timeout=120
            )
            r.raise_for_status()
            pred = r.json()
            if "id" not in pred: return f"❌ Replicate 回傳異常：{pred}"
            pred_id = pred["id"]
            pred_url = f"https://api.replicate.com/v1/predictions/{pred_id}"
            for _ in range(60):
                time.sleep(10)
                resp = requests.get(pred_url, headers=headers, timeout=15)
                resp.raise_for_status()
                data = resp.json()
                status = data.get("status")
                if status == "succeeded":
                    video_url = data.get("output")
                    if isinstance(video_url, list): video_url = video_url[0]
                    _download(video_url, out)
                    return f"✅ Replicate 影片已生成：{out}"
                elif status == "failed":
                    return f"❌ Replicate 生成失敗：{data.get('error', '未知錯誤')}"
            return "❌ Replicate 逾時（超過 10 分鐘）"

        elif provider == "runway":
            api_key = os.getenv("RUNWAY_API_KEY", "")
            if not api_key: return "❌ 缺少 RUNWAY_API_KEY，請在 .env 加入"
            try: import runwayml
            except ImportError: return "❌ 請先安裝：pip install runwayml"
            client = runwayml.RunwayML(api_key=api_key)
            if image_url:
                task = client.image_to_video.create(model="gen4_turbo", prompt_image=image_url,
                    prompt_text=prompt, duration=int(min(duration, 10)), ratio="1280:720")
            else:
                task = client.text_to_video.create(model="gen4_turbo", prompt_text=prompt,
                    duration=int(min(duration, 10)), ratio="1280:720")
            task_id = task.id
            for _ in range(60):
                time.sleep(5)
                t = client.tasks.retrieve(task_id)
                if t.status == "SUCCEEDED":
                    _download(t.output[0], out)
                    return f"✅ Runway 影片已生成：{out}"
                elif t.status in ("FAILED", "CANCELLED"):
                    return f"❌ Runway 生成失敗：{t.failure_reason}"
            return "❌ Runway 逾時"

        elif provider == "kling":
            import hmac, hashlib, base64, json as _json
            access_key = os.getenv("KLING_ACCESS_KEY", "")
            secret_key = os.getenv("KLING_SECRET_KEY", "")
            if not access_key or not secret_key:
                return "❌ 缺少 KLING_ACCESS_KEY / KLING_SECRET_KEY，請在 .env 加入"
            import jwt as _jwt
            payload_jwt = {"iss": access_key, "exp": int(time.time()) + 1800, "nbf": int(time.time()) - 5}
            token = _jwt.encode(payload_jwt, secret_key, algorithm="HS256")
            headers = {"Authorization": f"Bearer {token}", "Content-Type": "application/json"}
            body = {"model_name": "kling-v1", "prompt": prompt,
                    "duration": str(int(min(duration, 10))), "aspect_ratio": "16:9"}
            if image_url: body["image_url"] = image_url
            endpoint = "image2video" if image_url else "text2video"
            r = requests.post(f"https://api.klingai.com/v1/videos/{endpoint}",
                json=body, headers=headers, timeout=30)
            r.raise_for_status()
            data = r.json()
            if data.get("code") != 0: return f"❌ Kling API 錯誤：{data.get('message')}"
            task_id = data["data"]["task_id"]
            for _ in range(60):
                time.sleep(5)
                payload_jwt["exp"] = int(time.time()) + 1800
                token = _jwt.encode(payload_jwt, secret_key, algorithm="HS256")
                headers["Authorization"] = f"Bearer {token}"
                resp = requests.get(f"https://api.klingai.com/v1/videos/{endpoint}/{task_id}",
                    headers=headers, timeout=15)
                resp.raise_for_status()
                d = resp.json().get("data", {})
                status = d.get("task_status")
                if status == "succeed":
                    video_url = d["task_result"]["videos"][0]["url"]
                    _download(video_url, out)
                    return f"✅ Kling 影片已生成：{out}"
                elif status == "failed":
                    return f"❌ Kling 生成失敗：{d.get('task_status_msg')}"
            return "❌ Kling 逾時"
        else:
            return f"❌ 未知 provider：{provider}，可用：replicate / runway / kling"
    except Exception as e:
        return f"❌ AI 影片生成失敗：{e}\n{traceback.format_exc()}"


def execute_alert_monitor(action, name="", condition="", threshold="", target="", interval=30, chat_id=None, _bot_send=None):
    global _alert_monitors
    try:
        import threading, time, psutil
        if action == "list":
            if not _alert_monitors: return "⚠️ 無執行中的監控"
            return "📊 監控清單：\n" + "\n".join(f"- {k}: {v['condition']} {v['threshold']}" for k,v in _alert_monitors.items())
        elif action == "stop":
            if name in _alert_monitors: _alert_monitors[name]["running"] = False; del _alert_monitors[name]; return f"✅ 已停止監控：{name}"
            return f"⚠️ 找不到監控：{name}"
        elif action == "start":
            if name in _alert_monitors: return f"⚠️ 已有同名監控：{name}"
            cfg = {"condition": condition, "threshold": threshold, "target": target, "running": True}
            _alert_monitors[name] = cfg
            send_chat_id = chat_id or OWNER_ID
            def _monitor():
                import easyocr; reader = None
                while _alert_monitors.get(name, {}).get("running"):
                    try:
                        triggered = False; msg = ""; val = threshold
                        if condition == "cpu_above":
                            v = psutil.cpu_percent(1)
                            if v > float(val): triggered = True; msg = f"⚠️ CPU 使用率 {v:.1f}% 超過 {val}%"
                        elif condition == "mem_above":
                            v = psutil.virtual_memory().percent
                            if v > float(val): triggered = True; msg = f"⚠️ 記憶體使用率 {v:.1f}% 超過 {val}%"
                        elif condition == "disk_above":
                            v = psutil.disk_usage("/").percent
                            if v > float(val): triggered = True; msg = f"⚠️ 磁碟使用率 {v:.1f}% 超過 {val}%"
                        elif condition == "process_missing":
                            pnames = [p.name().lower() for p in psutil.process_iter(["name"])]
                            if not any(target.lower() in n for n in pnames): triggered = True; msg = f"⚠️ 程序 {target} 已停止執行"
                        elif condition == "process_running":
                            pnames = [p.name().lower() for p in psutil.process_iter(["name"])]
                            if any(target.lower() in n for n in pnames): triggered = True; msg = f"ℹ️ 程序 {target} 正在執行"
                        elif condition == "screen_text_found":
                            import pyautogui
                            if reader is None: reader = easyocr.Reader(["ch_tra","en"], gpu=False)
                            screenshot = pyautogui.screenshot()
                            import tempfile; tmp = tempfile.mktemp(suffix=".png"); screenshot.save(tmp)
                            results = reader.readtext(tmp, detail=0); full_text = " ".join(results)
                            Path(tmp).unlink(missing_ok=True)
                            if target.lower() in full_text.lower(): triggered = True; msg = f"ℹ️ 螢幕偵測到文字：{target}"
                        if triggered and _bot_send:
                            import asyncio
                            asyncio.run_coroutine_threadsafe(_bot_send(send_chat_id, f"🔔 【監控告警】{name}\n{msg}"), asyncio.get_event_loop())
                    except Exception: pass
                    time.sleep(int(interval))
            threading.Thread(target=_monitor, daemon=True).start()
            return f"✅ 監控已啟動：{name}（{condition} {threshold}，每 {interval}s 檢查）"
    except Exception as e:
        return f"❌ 監控告警失敗：{e}"


def execute_audio_process(action, input_path, output="", start_ms=0, end_ms=0):
    try:
        from pydub import AudioSegment
        if action == "convert":
            fmt = Path(output).suffix.lstrip(".")
            AudioSegment.from_file(input_path).export(output, format=fmt)
            return f"✅ 已轉換：{output}"
        elif action == "trim":
            audio = AudioSegment.from_file(input_path)[start_ms:end_ms]
            out = output or input_path.replace(".", "_trim.")
            audio.export(out, format=Path(out).suffix.lstrip("."))
            return f"✅ 已剪輯：{out}"
    except Exception as e:
        return f"❌ 音訊處理失敗：{e}"


def execute_audio_transcribe(action, path="", duration=30, language="", output=""):
    try:
        if action == "transcribe_file":
            if not path: return "❌ 需提供 path"
            try:
                import whisper; model = whisper.load_model("base"); result = model.transcribe(path, language=language or None); text = result["text"]
            except ImportError:
                import speech_recognition as _sr; r = _sr.Recognizer()
                with _sr.AudioFile(path) as src: audio = r.record(src)
                text = r.recognize_google(audio, language=language or "zh-TW")
            if output:
                with open(output, "w", encoding="utf-8") as f: f.write(text)
                return f"✅ 轉錄完成，已儲存：{output}\n\n{text[:500]}"
            return f"📝 轉錄結果：\n{text}"
        elif action == "transcribe_mic":
            import speech_recognition as _sr, time as _t
            r = _sr.Recognizer(); m = _sr.Microphone(); results = []; end = _t.time() + duration
            while _t.time() < end:
                with m as src:
                    r.adjust_for_ambient_noise(src, duration=0.3)
                    try: audio = r.listen(src, timeout=3, phrase_time_limit=10); text = r.recognize_google(audio, language=language or "zh-TW"); results.append(text)
                    except: pass
            full = " ".join(results)
            if output:
                with open(output, "w", encoding="utf-8") as f: f.write(full)
            return f"📝 麥克風轉錄：\n{full}" if full else "❌ 未偵測到語音"
        elif action == "transcribe_system": return "⚠️ 系統音訊轉錄需要虛擬音訊裝置（如 VB-Cable），建議先錄製後用 transcribe_file"
        return "未知動作"
    except ImportError: return "❌ 請先安裝：pip install SpeechRecognition（或 openai-whisper 提供更高精度）"
    except Exception as e: return f"❌ audio_transcribe 失敗：{e}"


def execute_auto_skill(action, goal="", skill_name="", code="", test_input=""):
    """自動技能生成與部署"""
    if action == "generate":
        try:
            from anthropic import Anthropic; c = Anthropic()
            resp = c.messages.create(model="claude-sonnet-4-6", max_tokens=2000,
                messages=[{"role": "user", "content": f"你是Python專家。請為小牛馬Telegram Bot生成一個新的execute_函數。\n技能需求：{goal}\n函數名稱：execute_{skill_name or 'new_skill'}\n要求：完整錯誤處理、回傳字串、使用標準庫。只輸出Python程式碼。"}])
            generated = resp.content[0].text
            draft = Path(__file__).parent / f"skill_{skill_name or 'draft'}.py"
            draft.write_text(generated, encoding="utf-8")
            return f"✅ 技能已生成\n儲存至：{draft}\n\n```python\n{generated[:600]}\n```"
        except Exception as e: return f"生成失敗：{e}"
    elif action == "test":
        try: compile(code, "<string>", "exec"); return "✅ 程式碼語法正確"
        except SyntaxError as e: return f"❌ 語法錯誤：{e}"
    elif action == "deploy":
        try:
            if not code:
                draft = Path(__file__).parent / f"skill_{skill_name}.py"
                if draft.exists(): code = draft.read_text(encoding="utf-8")
                else: return "請提供程式碼"
            bot_path = Path(__file__); content = bot_path.read_text(encoding="utf-8")
            marker = "\nasync def start("
            if marker in content:
                new_content = content.replace(marker, f"\n\n# ── 自動部署：{skill_name} ──\n{code}\n{marker}")
                bot_path.write_text(new_content, encoding="utf-8")
                return f"✅ 技能 {skill_name} 已部署，重啟後生效"
            return "部署失敗：找不到插入點"
        except Exception as e: return f"部署失敗：{e}"
    elif action == "list_skills":
        import re; content = Path(__file__).read_text(encoding="utf-8")
        skills = re.findall(r'def execute_(\w+)\(', content)
        return f"已部署技能：{len(skills)} 個\n" + "\n".join(f"• {s}" for s in sorted(set(skills)))
    return f"未知動作：{action}"


def execute_auto_trade(action, symbol="", amount=0.0, price=0.0, order_type="market", api_key="", api_secret=""):
    """加密貨幣自動交易：Binance API"""
    bk = api_key or os.environ.get("BINANCE_KEY", ""); bs = api_secret or os.environ.get("BINANCE_SECRET", "")
    if action == "price":
        try:
            sym = symbol.upper().replace("/", "").replace("-", "")
            resp = requests.get(f"https://api.binance.com/api/v3/ticker/24hr?symbol={sym}", timeout=5)
            if resp.status_code == 200:
                d = resp.json(); chg = float(d.get("priceChangePercent", "0")); arrow = "▲" if chg >= 0 else "▼"
                return f"{sym} 行情：\n價格：{d.get('lastPrice','')} USDT\n24h：{arrow} {chg}%\n高：{d.get('highPrice','')} | 低：{d.get('lowPrice','')} | 量：{d.get('volume','')}"
            return f"查詢失敗：{resp.status_code}"
        except Exception as e: return f"查詢失敗：{e}"
    elif action == "balance":
        if not bk or not bs: return "需要設定 BINANCE_KEY 和 BINANCE_SECRET"
        try:
            import hmac, hashlib; ts = int(time.time() * 1000); params = f"timestamp={ts}"
            sig = hmac.new(bs.encode(), params.encode(), hashlib.sha256).hexdigest()
            resp = requests.get(f"https://api.binance.com/api/v3/account?{params}&signature={sig}", headers={"X-MBX-APIKEY": bk}, timeout=10)
            if resp.status_code == 200:
                non_zero = [b for b in resp.json().get("balances", []) if float(b["free"]) > 0.0001]
                return "💰 帳戶餘額：\n" + "\n".join(f"  {b['asset']}: {b['free']}" for b in non_zero[:15])
            return f"查詢失敗：{resp.text}"
        except Exception as e: return f"查詢失敗：{e}"
    elif action in ("buy", "sell"):
        if not bk or not bs: return "需要設定 BINANCE_KEY 和 BINANCE_SECRET 才能下單"
        try:
            import hmac, hashlib; sym = symbol.upper().replace("/", "").replace("-", "")
            side = "BUY" if action == "buy" else "SELL"; ts = int(time.time() * 1000)
            params = f"symbol={sym}&side={side}&type={order_type.upper()}&quantity={amount}&timestamp={ts}"
            if order_type.lower() == "limit" and float(price) > 0: params += f"&price={price}&timeInForce=GTC"
            sig = hmac.new(bs.encode(), params.encode(), hashlib.sha256).hexdigest()
            resp = requests.post(f"https://api.binance.com/api/v3/order?{params}&signature={sig}", headers={"X-MBX-APIKEY": bk}, timeout=10)
            if resp.status_code == 200:
                d = resp.json(); return f"✅ 訂單成功\nID: {d.get('orderId','')}\n{side} {amount} {sym} @ {order_type}\n狀態：{d.get('status','')}"
            return f"下單失敗：{resp.text}"
        except Exception as e: return f"下單失敗：{e}"
    elif action == "open_orders":
        if not bk or not bs: return "需要設定 BINANCE_KEY 和 BINANCE_SECRET"
        try:
            import hmac, hashlib; ts = int(time.time() * 1000)
            sym_part = f"&symbol={symbol.upper()}" if symbol else ""; params = f"timestamp={ts}{sym_part}"
            sig = hmac.new(bs.encode(), params.encode(), hashlib.sha256).hexdigest()
            resp = requests.get(f"https://api.binance.com/api/v3/openOrders?{params}&signature={sig}", headers={"X-MBX-APIKEY": bk}, timeout=10)
            if resp.status_code == 200:
                orders = resp.json()
                if not orders: return "目前無掛單"
                return f"掛單列表（{len(orders)}筆）：\n" + "\n".join(f"  [{o['orderId']}] {o['side']} {o['origQty']} {o['symbol']} @ {o['price']}" for o in orders[:10])
            return f"查詢失敗：{resp.text}"
        except Exception as e: return f"查詢失敗：{e}"
    return f"未知動作：{action}"


def execute_automation(action, condition_type="", condition_value="", command="",
                       duration=60.0, layout="side_by_side", x=0, y=0, w=0, h=0,
                       keyword="", output=""):
    try:
        if action == "if_then":
            import subprocess, psutil, time, os
            deadline = time.time() + float(duration)
            while time.time() < deadline:
                triggered = False
                if condition_type == "cpu_above": triggered = psutil.cpu_percent(1) > float(condition_value)
                elif condition_type == "mem_above": triggered = psutil.virtual_memory().percent > float(condition_value)
                elif condition_type == "file_exists": triggered = Path(condition_value).exists()
                elif condition_type == "process_running": triggered = any(condition_value.lower() in p.name().lower() for p in psutil.process_iter())
                elif condition_type == "time_is": triggered = dt.dt.datetime.now().strftime("%H:%M") == condition_value
                if triggered:
                    subprocess.Popen(command, shell=True)
                    return f"✅ 條件達成（{condition_type}={condition_value}），已執行：{command}"
                time.sleep(2)
            return f"⏳ 監控 {duration}s 內條件未達成"
        elif action == "window_arrange":
            import win32gui, win32con, ctypes
            user32 = ctypes.windll.user32; sw = user32.GetSystemMetrics(0); sh = user32.GetSystemMetrics(1)
            hwnds = []
            win32gui.EnumWindows(lambda h, l: l.append(h) if win32gui.IsWindowVisible(h) and win32gui.GetWindowText(h) else None, hwnds)
            if layout == "side_by_side" and len(hwnds) >= 2:
                win32gui.MoveWindow(hwnds[0], 0, 0, sw//2, sh, True)
                win32gui.MoveWindow(hwnds[1], sw//2, 0, sw//2, sh, True)
                return f"✅ 左右排列完成（{sw}x{sh}）"
            elif layout == "quad" and len(hwnds) >= 4:
                positions = [(0,0,sw//2,sh//2),(sw//2,0,sw//2,sh//2),(0,sh//2,sw//2,sh//2),(sw//2,sh//2,sw//2,sh//2)]
                for i, (x_,y_,w_,h_) in enumerate(positions[:4]):
                    win32gui.MoveWindow(hwnds[i], x_, y_, w_, h_, True)
                return "✅ 四象限排列完成"
            elif layout == "stack":
                h_each = sh // max(len(hwnds), 1)
                for i, hwnd in enumerate(hwnds[:8]):
                    win32gui.MoveWindow(hwnd, 0, i*h_each, sw, h_each, True)
                return f"✅ 堆疊排列完成（{len(hwnds[:8])} 個視窗）"
            elif layout == "maximize_all":
                for hwnd in hwnds: win32gui.ShowWindow(hwnd, win32con.SW_MAXIMIZE)
                return f"✅ 全部最大化（{len(hwnds)} 個）"
            return f"✅ 排列完成"
        elif action == "region_ocr":
            import pyautogui, easyocr
            region = (int(x), int(y), int(w), int(h)) if w and h else None
            screenshot = pyautogui.screenshot(region=region)
            out = output or str(Path("C:/Users/blue_/Desktop/測試檔案") / "region_ocr.png")
            screenshot.save(out)
            reader = easyocr.Reader(["ch_tra","en"], gpu=False)
            results = reader.readtext(out, detail=0)
            return f"🔍 區域 OCR 結果：\n" + "\n".join(results)
        elif action == "window_screenshot":
            import win32gui, pyautogui
            hwnds = []
            win32gui.EnumWindows(lambda h, l: l.append(h) if win32gui.IsWindowVisible(h) and keyword.lower() in win32gui.GetWindowText(h).lower() else None, hwnds)
            if not hwnds: return f"⚠️ 找不到視窗：{keyword}"
            hwnd = hwnds[0]; win32gui.SetForegroundWindow(hwnd)
            rect = win32gui.GetWindowRect(hwnd)
            x_, y_, x2, y2 = rect; w_, h_ = x2-x_, y2-y_
            screenshot = pyautogui.screenshot(region=(x_, y_, w_, h_))
            out = output or str(Path("C:/Users/blue_/Desktop/測試檔案") / f"window_{keyword}_{dt.dt.datetime.now().strftime('%H%M%S')}.png")
            screenshot.save(out)
            return f"✅ 視窗截圖：{out}"
    except Exception as e:
        return f"❌ 自動化失敗：{e}"


def execute_barcode(image_path=""):
    try:
        from pyzbar.pyzbar import decode; from PIL import Image
        img = Image.open(image_path) if image_path else pyautogui.screenshot(); results = decode(img)
        if not results: return "❌ 未偵測到條碼"
        return "\n".join(f"類型：{r.type}  內容：{r.data.decode('utf-8', errors='replace')}" for r in results)
    except Exception as e: return f"❌ 條碼掃描失敗：{e}"


def execute_bluetooth(action, mac=""):
    try:
        import asyncio, bleak
        if action == "scan":
            async def _scan():
                return await bleak.BleakScanner.discover(timeout=5.0)
            devices = asyncio.run(_scan())
            return "\n".join(f"{d.address} {d.name or '(未知)'}" for d in devices) or "找不到裝置"
        elif action == "connect":
            async def _conn():
                async with bleak.BleakClient(mac) as c:
                    return f"已連線：{mac}（服務數：{len(c.services)}）"
            return asyncio.run(_conn())
    except Exception as e:
        return f"藍牙操作失敗：{e}"


def execute_browser_advanced(action, selector="", value="", name="", tab_index=0, timeout=30.0, url_pattern=""):
    global _browser_page
    try:
        from playwright.sync_api import sync_playwright
        import json
        if _browser_page is None or _browser_page.is_closed():
            return "⚠️ 瀏覽器未開啟，請先用 browser_control 開啟瀏覽器"
        page = _browser_page
        if action == "wait_element": page.wait_for_selector(selector, timeout=float(timeout)*1000); return f"✅ 元素已出現：{selector}"
        elif action == "get_cookies":
            cookies = page.context.cookies()
            lines = [f"{c['name']}={c['value'][:30]}" for c in cookies[:20]]
            return "🍪 Cookies：\n" + "\n".join(lines)
        elif action == "set_cookie": page.context.add_cookies([{"name": name, "value": value, "url": page.url}]); return f"✅ Cookie 已設定：{name}={value}"
        elif action == "list_tabs":
            pages = page.context.pages
            lines = [f"{i}: {p.title()} - {p.url[:50]}" for i, p in enumerate(pages)]
            return "📑 所有分頁：\n" + "\n".join(lines)
        elif action == "switch_tab":
            pages = page.context.pages
            if tab_index < len(pages): pages[tab_index].bring_to_front(); _browser_page = pages[tab_index]; return f"✅ 已切換到分頁 {tab_index}"
            return f"⚠️ 分頁 {tab_index} 不存在"
        elif action == "new_tab":
            new_page = page.context.new_page()
            if value: new_page.goto(value)
            _browser_page = new_page; return f"✅ 已開新分頁：{value or '空白'}"
        elif action == "close_tab":
            page.close(); pages = page.context.pages
            if pages: _browser_page = pages[-1]
            return "✅ 已關閉目前分頁"
        elif action == "fill_form":
            import json as _json
            fields = _json.loads(value) if value.startswith("{") else {}
            for sel, val in fields.items(): page.fill(sel, val)
            return f"✅ 表單已填寫"
        elif action == "select_option": page.select_option(selector, value); return f"✅ 已選擇：{value}"
        elif action == "scroll_to": page.locator(selector).scroll_into_view_if_needed(); return f"✅ 已滾動到：{selector}"
        elif action == "get_html": return page.inner_html(selector or "body")[:2000]
        elif action == "wait_url": page.wait_for_url(f"**{url_pattern}**", timeout=float(timeout)*1000); return f"✅ URL 已包含：{url_pattern}"
    except Exception as e:
        return f"❌ 瀏覽器進階操作失敗：{e}"


def execute_browser_control(action: str, url: str = "", selector: str = "", text: str = "") -> str:
    try:
        from playwright.sync_api import sync_playwright
    except (ImportError, OSError) as _dll_err:
        if action in ("open", "goto") and url:
            import subprocess as _sp; _sp.Popen(f'start "" "{url}"', shell=True)
            return f"已用瀏覽器開啟：{url}"
        return f"瀏覽器功能暫不可用（DLL 問題）：{_dll_err}"
    try:
        if action == "open":
            if _browser_ctx.get("page"): _browser_ctx["browser"].close(); _browser_ctx["pw"].stop(); _browser_ctx.clear()
            pw = sync_playwright().start(); b = pw.chromium.launch(headless=True)
            page = b.new_page(); page.goto(url or "https://www.google.com")
            _browser_ctx.update({"pw": pw, "browser": b, "page": page})
            return f"已開啟：{url}"
        page = _browser_ctx.get("page")
        if not page: return "瀏覽器未開啟，請先使用 open 開啟網頁"
        if action == "goto": page.goto(url); return f"已前往：{url}"
        elif action == "click": page.click(selector); return f"已點擊：{selector}"
        elif action == "type": page.fill(selector, text); return f"已輸入到 {selector}：{text}"
        elif action == "get_text": return page.inner_text(selector or "body")[:2000]
        elif action == "screenshot":
            img_bytes = page.screenshot()
            return f"__BROWSER_SCREENSHOT__:{img_bytes.hex()}"
        elif action == "close":
            _browser_ctx["browser"].close(); _browser_ctx["pw"].stop(); _browser_ctx.clear()
            return "瀏覽器已關閉"
        return f"未知動作：{action}"
    except Exception as e:
        return f"執行失敗：{str(e)}"


def execute_calendar(action, days=7, title="", start="", end="", description=""):
    try:
        from google.oauth2.credentials import Credentials; from googleapiclient.discovery import build
        creds_path = Path("C:/Users/blue_/claude-telegram-bot/gcal_token.json")
        if not creds_path.exists(): return "❌ 未找到 Google Calendar 憑證"
        creds = Credentials.from_authorized_user_file(str(creds_path)); service = build("calendar", "v3", credentials=creds)
        if action == "list":
            now = dt.datetime.now(dt.timezone.utc)
            events = service.events().list(calendarId="primary", timeMin=now.isoformat(), timeMax=(now + dt.timedelta(days=days)).isoformat(), maxResults=20, singleEvents=True, orderBy="startTime").execute().get("items", [])
            if not events: return f"未來 {days} 天沒有行程"
            return "\n".join(f"📅 {e['start'].get('dateTime',e['start'].get('date'))}  {e.get('summary','（無標題）')}" for e in events)
        elif action == "add":
            event = {"summary": title, "description": description, "start": {"dateTime": start, "timeZone": "Asia/Taipei"}, "end": {"dateTime": end, "timeZone": "Asia/Taipei"}}
            created = service.events().insert(calendarId="primary", body=event).execute()
            return f"✅ 行程已新增：{created.get('summary')}"
        elif action == "delete":
            if not title: return "❌ 請提供要刪除的行程標題"
            events = service.events().list(calendarId="primary", q=title, maxResults=5, singleEvents=True).execute().get("items", [])
            if not events: return f"❌ 找不到符合的行程：{title}"
            service.events().delete(calendarId="primary", eventId=events[0]["id"]).execute()
            return f"✅ 已刪除行程：{events[0].get('summary', title)}"
    except Exception as e: return f"❌ 行事曆操作失敗：{e}"


def execute_chart(chart_type, data_json, title="", output=""):
    try:
        import matplotlib; matplotlib.use("Agg"); import matplotlib.pyplot as plt, json
        data = json.loads(data_json)
        out_path = output or str(Path("C:/Users/blue_/Desktop/測試檔案") / f"chart_{dt.dt.datetime.now().strftime('%H%M%S')}.png")
        fig, ax = plt.subplots()
        if chart_type == "line":
            for label, values in data.items(): ax.plot(values, label=label)
            ax.legend()
        elif chart_type == "bar": ax.bar(list(data.keys()), list(data.values()))
        elif chart_type == "pie": ax.pie(list(data.values()), labels=list(data.keys()), autopct="%1.1f%%")
        if title: ax.set_title(title)
        plt.tight_layout(); plt.savefig(out_path); plt.close()
        return out_path
    except Exception as e: return f"❌ 圖表生成失敗：{e}"


def execute_clipboard(action, text=""):
    import pyperclip
    if action == "get":
        return pyperclip.paste() or "（剪貼簿是空的）"
    else:
        pyperclip.copy(text)
        return f"已寫入剪貼簿：{text}"


def execute_clipboard_history(action, index=0):
    global _clipboard_hist
    try:
        import pyperclip, threading, time
        if action == "start_watch":
            def _watch():
                last = ""
                while True:
                    try:
                        cur = pyperclip.paste()
                        if cur != last and cur: _clipboard_hist.insert(0, cur); last = cur
                        if len(_clipboard_hist) > 50: _clipboard_hist.pop()
                    except Exception: pass
                    time.sleep(1)
            threading.Thread(target=_watch, daemon=True).start()
            return "✅ 剪貼簿歷史監控已啟動"
        elif action == "list":
            if not _clipboard_hist: return "⚠️ 剪貼簿歷史為空（請先執行 start_watch）"
            return "📋 剪貼簿歷史：\n" + "\n".join(f"[{i}] {item[:80]}" for i, item in enumerate(_clipboard_hist[:20]))
        elif action == "get":
            if index < len(_clipboard_hist): return f"📋 [{index}]：{_clipboard_hist[index]}"
            return f"⚠️ 索引 {index} 超出範圍"
        elif action == "set":
            if index < len(_clipboard_hist): pyperclip.copy(_clipboard_hist[index]); return f"✅ 已復原剪貼簿 [{index}]"
            return f"⚠️ 索引 {index} 超出範圍"
        elif action == "clear": _clipboard_hist.clear(); return "✅ 剪貼簿歷史已清除"
    except Exception as e:
        return f"❌ 剪貼簿歷史失敗：{e}"


def execute_clipboard_image(action, path=""):
    try:
        import win32clipboard; from PIL import Image; import io as _io
        if action == "get":
            win32clipboard.OpenClipboard()
            try: data = win32clipboard.GetClipboardData(win32clipboard.CF_DIB)
            finally: win32clipboard.CloseClipboard()
            out = path or str(Path("C:/Users/blue_/Desktop/測試檔案") / f"clipboard_{dt.dt.datetime.now().strftime('%H%M%S')}.png")
            Image.open(_io.BytesIO(data)).save(out); return f"✅ 剪貼簿圖片已存：{out}"
        elif action == "set":
            img = Image.open(path).convert("RGB"); buf = _io.BytesIO(); img.save(buf,"BMP"); data = buf.getvalue()[14:]
            win32clipboard.OpenClipboard()
            try: win32clipboard.EmptyClipboard(); win32clipboard.SetClipboardData(win32clipboard.CF_DIB, data)
            finally: win32clipboard.CloseClipboard()
            return f"✅ 圖片已複製到剪貼簿"
    except Exception as e:
        return f"❌ 剪貼簿圖片失敗：{e}"


def execute_cloud_storage(action, path, drive_id="root"):
    try:
        from google.oauth2.credentials import Credentials
        from googleapiclient.discovery import build
        from googleapiclient.http import MediaFileUpload, MediaIoBaseDownload
        import io as _io
        creds_path = Path("C:/Users/blue_/claude-telegram-bot/gdrive_token.json")
        if not creds_path.exists():
            return "❌ 未找到 Google Drive 憑證（gdrive_token.json）"
        creds = Credentials.from_authorized_user_file(str(creds_path))
        service = build("drive", "v3", credentials=creds)
        if action == "upload":
            media = MediaFileUpload(path)
            meta = {"name": Path(path).name, "parents": [drive_id]}
            f = service.files().create(body=meta, media_body=media, fields="id").execute()
            return f"✅ 已上傳，檔案 ID：{f.get('id')}"
        elif action == "download":
            req = service.files().get_media(fileId=drive_id)
            buf = _io.BytesIO()
            dl = MediaIoBaseDownload(buf, req)
            done = False
            while not done:
                _, done = dl.next_chunk()
            with open(path, "wb") as f:
                f.write(buf.getvalue())
            return f"✅ 已下載到：{path}"
        return "未知動作"
    except Exception as e:
        return f"雲端儲存失敗：{e}"


def execute_color_pick(action, x=0, y=0, region_w=100, region_h=100):
    try:
        import pyautogui; from PIL import Image
        if action == "pick":
            screenshot = pyautogui.screenshot(); pixel = screenshot.getpixel((int(x), int(y))); r, g, b = pixel[:3]
            return f"🎨 座標 ({x},{y}) 的顏色：\nRGB: ({r}, {g}, {b})\nHEX: #{r:02X}{g:02X}{b:02X}"
        elif action == "dominant":
            screenshot = pyautogui.screenshot(region=(int(x), int(y), int(region_w), int(region_h)))
            img = screenshot.convert("RGB").resize((50, 50)); pixels = list(img.getdata())
            from collections import Counter; most_common = Counter(pixels).most_common(5)
            lines = [f"RGB({r},{g},{b}) = #{r:02X}{g:02X}{b:02X}  出現 {cnt} 次" for (r,g,b), cnt in most_common]
            return f"🎨 區域主要顏色：\n" + "\n".join(lines)
    except Exception as e: return f"❌ 顏色選取失敗：{e}"


def execute_com_auto(app, action, path="", sheet=None, cell="", value="", macro="", to="", subject=""):
    try:
        import win32com.client as _com
        if app == "excel":
            xl = _com.Dispatch("Excel.Application"); xl.Visible = False
            wb = xl.Workbooks.Open(path) if path else (xl.Workbooks(1) if xl.Workbooks.Count > 0 else xl.Workbooks.Add())
            ws = wb.Sheets(sheet) if sheet else wb.ActiveSheet
            if action == "read_cell": return f"{cell} = {ws.Range(cell).Value}"
            elif action == "write_cell": ws.Range(cell).Value = value; return f"✅ 已寫入 {cell} = {value}"
            elif action == "run_macro": xl.Run(macro); return f"✅ 巨集 {macro} 已執行"
            elif action == "save": wb.Save(); return "✅ 已儲存"
            elif action == "list_sheets": return "\n".join(wb.Sheets(i+1).Name for i in range(wb.Sheets.Count))
            elif action == "close": wb.Close(SaveChanges=False); return "✅ 已關閉"
        elif app == "word":
            wd = _com.Dispatch("Word.Application"); wd.Visible = False
            doc = wd.Documents.Open(path) if path else wd.Documents.Add()
            if action == "read": return doc.Content.Text[:1000]
            elif action == "write": doc.Content.Text = value; return "✅ 已寫入"
            elif action == "save": doc.Save(); return "✅ 已儲存"
            elif action == "close": doc.Close(SaveChanges=False); return "✅ 已關閉"
        elif app == "outlook":
            ol = _com.Dispatch("Outlook.Application")
            if action == "send":
                mail = ol.CreateItem(0); mail.To = to; mail.Subject = subject; mail.Body = value; mail.Send()
                return f"✅ 郵件已發送至 {to}"
            elif action == "list_inbox":
                inbox = ol.GetNamespace("MAPI").GetDefaultFolder(6); items = inbox.Items; items.Sort("[ReceivedTime]", True)
                result = [f"寄件人：{item.SenderName} | 主旨：{item.Subject}" for i, item in enumerate(items) if i < 10]
                return "\n".join(result)
        return f"✅ {app} {action} 完成"
    except Exception as e:
        return f"❌ com_auto 失敗：{e}"


def execute_data_process(action, path="", output="", query="", data="", paths=""):
    try:
        import json, csv, io as _io
        if action == "read_json":
            content = json.loads(Path(path).read_text(encoding="utf-8"))
            if isinstance(content, list):
                return f"📄 JSON（{len(content)} 筆）：\n" + json.dumps(content[:5], ensure_ascii=False, indent=2)
            return f"📄 JSON：\n" + json.dumps(content, ensure_ascii=False, indent=2)[:2000]
        elif action == "write_json":
            obj = json.loads(data)
            Path(output or path).write_text(json.dumps(obj, ensure_ascii=False, indent=2), encoding="utf-8")
            return f"✅ JSON 已儲存：{output or path}"
        elif action == "read_csv":
            rows = []
            with open(path, encoding="utf-8-sig", errors="replace") as f:
                reader = csv.DictReader(f); rows = list(reader)
            return f"📊 CSV（{len(rows)} 筆，欄位：{list(rows[0].keys()) if rows else []}）：\n" + json.dumps(rows[:5], ensure_ascii=False, indent=2)
        elif action == "write_csv":
            obj = json.loads(data)
            if not obj: return "⚠️ 資料為空"
            out = output or path
            with open(out, "w", newline="", encoding="utf-8-sig") as f:
                writer = csv.DictWriter(f, fieldnames=obj[0].keys()); writer.writeheader(); writer.writerows(obj)
            return f"✅ CSV 已儲存：{out}"
        elif action == "filter":
            with open(path, encoding="utf-8-sig", errors="replace") as f: rows = list(csv.DictReader(f))
            filtered = []
            for row in rows:
                try:
                    if eval(query, {"__builtins__": {}}, row): filtered.append(row)
                except Exception: pass
            return f"📊 過濾結果（{len(filtered)}/{len(rows)} 筆）：\n" + json.dumps(filtered[:10], ensure_ascii=False, indent=2)
        elif action == "stats":
            with open(path, encoding="utf-8-sig", errors="replace") as f: rows = list(csv.DictReader(f))
            if not rows: return "⚠️ 無資料"
            stats = {}
            for col in rows[0].keys():
                vals = [row[col] for row in rows if row[col]]
                try:
                    nums = [float(v) for v in vals]
                    stats[col] = {"count": len(nums), "min": min(nums), "max": max(nums), "avg": round(sum(nums)/len(nums),2)}
                except Exception:
                    stats[col] = {"count": len(vals), "unique": len(set(vals))}
            return "📊 統計：\n" + json.dumps(stats, ensure_ascii=False, indent=2)
        elif action == "convert":
            ext_in = Path(path).suffix.lower(); ext_out = Path(output).suffix.lower() if output else ""
            if ext_in == ".json" and ext_out == ".csv":
                obj = json.loads(Path(path).read_text(encoding="utf-8"))
                if not isinstance(obj, list): obj = [obj]
                with open(output, "w", newline="", encoding="utf-8-sig") as f:
                    w = csv.DictWriter(f, fieldnames=obj[0].keys()); w.writeheader(); w.writerows(obj)
                return f"✅ JSON → CSV：{output}"
            elif ext_in == ".csv" and ext_out == ".json":
                with open(path, encoding="utf-8-sig", errors="replace") as f: rows = list(csv.DictReader(f))
                Path(output).write_text(json.dumps(rows, ensure_ascii=False, indent=2), encoding="utf-8")
                return f"✅ CSV → JSON：{output}"
            return f"⚠️ 不支援轉換：{ext_in} → {ext_out}"
        elif action == "merge":
            all_rows = []
            for p in paths.split(","):
                p = p.strip()
                if p.endswith(".csv"):
                    with open(p, encoding="utf-8-sig", errors="replace") as f: all_rows.extend(list(csv.DictReader(f)))
                elif p.endswith(".json"):
                    obj = json.loads(Path(p).read_text(encoding="utf-8"))
                    if isinstance(obj, list): all_rows.extend(obj)
            out = output or str(Path(paths.split(",")[0].strip()).parent / "merged.csv")
            with open(out, "w", newline="", encoding="utf-8-sig") as f:
                w = csv.DictWriter(f, fieldnames=all_rows[0].keys()); w.writeheader(); w.writerows(all_rows)
            return f"✅ 已合併 {len(all_rows)} 筆資料 → {out}"
        elif action == "to_table":
            with open(path, encoding="utf-8-sig", errors="replace") as f: rows = list(csv.DictReader(f))
            if not rows: return "⚠️ 無資料"
            cols = list(rows[0].keys())
            header = " | ".join(cols); sep = "-" * len(header)
            lines = [header, sep] + [" | ".join(str(r.get(c,"")) for c in cols) for r in rows[:20]]
            return "\n".join(lines)
    except Exception as e:
        return f"❌ 資料處理失敗：{e}"


def execute_database(type_, db, sql, name=""):
    try:
        if type_ == "sqlite":
            conn = sqlite3.connect(db); cur = conn.cursor(); cur.execute(sql); rows = cur.fetchall(); conn.commit(); conn.close()
            if rows: return "\n".join(str(r) for r in rows[:20])
            return "✅ 執行成功（無回傳資料）"
        elif type_ == "mysql":
            import pymysql; conn = pymysql.connect(host=db, database=name, read_default_file="~/.my.cnf")
            cur = conn.cursor(); cur.execute(sql); rows = cur.fetchall(); conn.commit(); conn.close()
            if rows: return "\n".join(str(r) for r in rows[:20])
            return "✅ 執行成功（無回傳資料）"
        return "未知類型"
    except Exception as e: return f"資料庫操作失敗：{e}"


def execute_datetime_config(action, timezone="", datetime_str=""):
    try:
        if action == "get":
            r = subprocess.run(["powershell", "-Command", "Get-Date | Format-List; (Get-TimeZone).DisplayName"], capture_output=True, text=True, encoding="utf-8", errors="replace")
            return f"🕐 系統時間：\n{r.stdout.strip()}"
        elif action == "sync": subprocess.run(["powershell", "-Command", "Start-Service w32tm -ErrorAction SilentlyContinue; w32tm /resync /force"], capture_output=True); return "✅ 已同步網路時間"
        elif action == "set_timezone":
            r = subprocess.run(["powershell", "-Command", f"Set-TimeZone -Id '{timezone}'"], capture_output=True, text=True)
            return f"✅ 時區已設定為：{timezone}" if r.returncode == 0 else f"❌ 設定失敗：{r.stderr.strip()}"
        elif action == "set_time":
            r = subprocess.run(["powershell", "-Command", f"Set-Date -Date '{datetime_str}'"], capture_output=True, text=True)
            return f"✅ 系統時間已設定為：{datetime_str}" if r.returncode == 0 else f"❌ 設定失敗：{r.stderr.strip()}"
    except Exception as e: return f"❌ 時間設定失敗：{e}"


def execute_ddg_search(query: str, region: str = "zh-tw", max_results: int = 5) -> str:
    try:
        from ddgs import DDGS
        max_results = min(max(max_results, 1), 10)
        results = []
        with DDGS() as ddgs:
            for r in ddgs.text(query, region=region, max_results=max_results):
                title = r.get("title", "")
                body = r.get("body", "")
                href = r.get("href", "")
                # 跳過大陸域名（.cn）避免超時
                if ".cn/" in href and "taiwan" not in href.lower():
                    pass
                results.append(f"🔍 {title}\n   {body[:120]}\n   {href}")
        return "\n\n".join(results) if results else "無搜尋結果"
    except Exception as e:
        return f"DuckDuckGo 搜尋失敗：{e}"


def execute_desktop_control(action: str, x=None, y=None, text=None, app=None, direction="down", amount=3, monitor=None) -> dict:
    """執行桌面控制動作，支援多螢幕，回傳 {"ok": bool, "message": str, "screenshot": bytes or None}"""
    global pyautogui
    if pyautogui is None:
        import pyautogui as _pag
        _pag.FAILSAFE = True
        pyautogui = _pag
    try:
        screenshot_bytes = None

        if action == "list_monitors":
            import mss
            with mss.mss() as sct:
                monitors = sct.monitors[1:]
            lines = []
            for i, m in enumerate(monitors, 1):
                lines.append(f"螢幕{i}：左={m['left']} 上={m['top']} 寬={m['width']} 高={m['height']}")
            return {"ok": True, "message": "\n".join(lines), "screenshot": None}

        elif action == "screenshot":
            from PIL import Image as _PIL_Image
            # 螢幕 mapping（由測試確認）：
            # 螢幕1 → dxcam output 0
            # 螢幕2 → 獨立顯示輸出，用 DPI-unaware GDI BitBlt
            # 螢幕3 → dxcam output 1
            _DXCAM_MAP = {1: 0, 3: 1}  # Windows 螢幕編號 → dxcam output_idx

            if monitor and monitor in _DXCAM_MAP:
                try:
                    import dxcam as _dxcam
                    _cam = _dxcam.create(output_idx=_DXCAM_MAP[monitor])
                    _frame = _cam.grab()
                    del _cam
                    if _frame is None:
                        raise RuntimeError("grab() 回傳 None")
                    img = _PIL_Image.fromarray(_frame)
                    label = f"螢幕{monitor}"
                except Exception as _e:
                    return {"ok": False, "message": f"截圖失敗：{_e}", "screenshot": None}

            elif monitor == 2:
                # 螢幕2 接在不同輸出，用 DPI-unaware GDI BitBlt
                try:
                    import ctypes as _ct, win32gui as _w32g, win32ui as _w32u, win32con as _w32c, mss as _mss
                    with _mss.mss() as sct:
                        _m2 = sct.monitors[2]  # mss monitors[2] = Telegram 螢幕（筆電left=3840、桌機left=1920）
                    _left, _top = _m2["left"], _m2["top"]
                    _w, _h = _m2["width"], _m2["height"]
                    _user32 = _ct.windll.user32
                    _old_ctx = _user32.SetThreadDpiAwarenessContext(_ct.c_void_p(-1))  # DPI_AWARENESS_CONTEXT_UNAWARE
                    try:
                        _hdesk = _w32g.GetDesktopWindow()
                        _hwdc = _w32g.GetWindowDC(_hdesk)
                        _mdc = _w32u.CreateDCFromHandle(_hwdc)
                        _sdc = _mdc.CreateCompatibleDC()
                        _bmp = _w32u.CreateBitmap()
                        _bmp.CreateCompatibleBitmap(_mdc, _w, _h)
                        _sdc.SelectObject(_bmp)
                        _sdc.BitBlt((0, 0), (_w, _h), _mdc, (_left, _top), _w32c.SRCCOPY)
                        _info = _bmp.GetInfo()
                        _bits = _bmp.GetBitmapBits(True)
                        img = _PIL_Image.frombuffer("RGB", (_info["bmWidth"], _info["bmHeight"]), _bits, "raw", "BGRX", 0, 1)
                        _w32g.DeleteObject(_bmp.GetHandle())
                        _sdc.DeleteDC()
                        _mdc.DeleteDC()
                        _w32g.ReleaseDC(_hdesk, _hwdc)
                    finally:
                        _user32.SetThreadDpiAwarenessContext(_old_ctx)
                    label = "螢幕2"
                except Exception as _e:
                    return {"ok": False, "message": f"螢幕2截圖失敗：{_e}", "screenshot": None}

            else:
                # 全螢幕：用 mss
                import mss as _mss
                with _mss.mss() as sct:
                    sct_img = sct.grab(sct.monitors[0])
                img = _PIL_Image.frombytes("RGB", sct_img.size, sct_img.rgb)
                label = "全螢幕"

            buf = io.BytesIO()
            img.save(buf, format="PNG")
            screenshot_bytes = buf.getvalue()
            return {"ok": True, "message": f"{label}截圖完成", "screenshot": screenshot_bytes}

        elif action == "click":
            ax, ay = _resolve_coords(x, y, monitor)
            pyautogui.click(ax, ay)
            return {"ok": True, "message": f"已點擊 ({x}, {y})" + (f" 螢幕{monitor}" if monitor else ""), "screenshot": None}

        elif action == "double_click":
            ax, ay = _resolve_coords(x, y, monitor)
            pyautogui.doubleClick(ax, ay)
            return {"ok": True, "message": f"已雙擊 ({x}, {y})" + (f" 螢幕{monitor}" if monitor else ""), "screenshot": None}

        elif action == "right_click":
            ax, ay = _resolve_coords(x, y, monitor)
            pyautogui.rightClick(ax, ay)
            return {"ok": True, "message": f"已右鍵點擊 ({x}, {y})" + (f" 螢幕{monitor}" if monitor else ""), "screenshot": None}

        elif action == "move":
            ax, ay = _resolve_coords(x, y, monitor)
            pyautogui.moveTo(ax, ay, duration=0.3)
            return {"ok": True, "message": f"滑鼠已移動到 ({x}, {y})" + (f" 螢幕{monitor}" if monitor else ""), "screenshot": None}

        elif action == "type":
            global _last_opened_hwnd
            import time as _t, ctypes as _ct
            _text_str = str(text)
            _injected = False

            # Win32 直接注射：優先用 open_app 記下的 HWND，否則搜尋標題
            try:
                import win32gui as _w32g, win32con as _w32c
                # 候選視窗：先放 open_app 記下的那個
                _parent_hwnds = []
                if _last_opened_hwnd and _w32g.IsWindow(_last_opened_hwnd):
                    _parent_hwnds.append(_last_opened_hwnd)
                # 再搜標題作備用
                _KEYS = ['記事本', 'notepad', 'untitled', '無標題', 'wordpad']
                def _enum_wins(hwnd, _):
                    if hwnd not in _parent_hwnds and _w32g.IsWindowVisible(hwnd):
                        t = _w32g.GetWindowText(hwnd).lower()
                        if any(k in t for k in _KEYS):
                            _parent_hwnds.append(hwnd)
                    return True
                _w32g.EnumWindows(_enum_wins, None)
                for _ph in _parent_hwnds:
                    _edit = [None]
                    def _find_edit_recursive(hwnd, _):
                        cn = _w32g.GetClassName(hwnd)
                        if cn in ('RichEditD2DPT', 'Edit', 'RICHEDIT50W'):
                            _edit[0] = hwnd
                            return False  # 停止枚舉
                        return True
                    try:
                        _w32g.EnumChildWindows(_ph, _find_edit_recursive, None)
                    except Exception:
                        pass
                    if _edit[0]:
                        _w32g.SendMessage(_edit[0], _w32c.EM_SETSEL, -1, -1)
                        _r = _ct.windll.user32.SendMessageW(_edit[0], _w32c.EM_REPLACESEL, True, _text_str)
                        _injected = True
                        break
            except Exception:
                pass

            if not _injected:
                import pyperclip
                pyperclip.copy(_text_str)
                _t.sleep(0.3)
                pyautogui.hotkey("ctrl", "v")

            _t.sleep(0.1)
            return {"ok": True, "message": f"已輸入文字：{text}", "screenshot": None}

        elif action == "press_key":
            pyautogui.press(text)
            return {"ok": True, "message": f"已按下按鍵：{text}", "screenshot": None}

        elif action == "open_app":
            import time as _t2, win32gui as _w32g_oa
            # 常見 App 名稱 → 實際啟動指令
            _app_alias = {
                "google chrome": "start chrome", "chrome": "start chrome",
                "firefox": "start firefox", "edge": "start msedge", "microsoft edge": "start msedge",
                "notepad": "start notepad", "記事本": "start notepad",
                "calculator": "start calc", "計算機": "start calc",
                "explorer": "start explorer", "檔案總管": "start explorer",
                "cmd": "start cmd", "命令提示字元": "start cmd",
                "powershell": "start powershell",
                "word": "start winword", "excel": "start excel", "powerpoint": "start powerpnt",
                "spotify": "start spotify", "discord": "start discord",
                "vscode": "start code", "visual studio code": "start code",
                "telegram": "start telegram",
                "line": "start LINE",
                "youtube": "start https://www.youtube.com",
                "google": "start https://www.google.com",
            }
            _cmd = _app_alias.get(app.lower().strip(), None)
            if _cmd is None:
                # URL → 用 webbrowser 模組開（比 start 更可靠）
                if app.strip().startswith("http://") or app.strip().startswith("https://"):
                    import webbrowser
                    webbrowser.open(app.strip())
                    import time as _t_url
                    _t_url.sleep(1.5)
                    return {"ok": True, "message": f"已在瀏覽器開啟：{app.strip()}", "screenshot": None}
                _cmd = f"start \"\" \"{app}\""
            # 記下開啟前已有的視窗 HWND
            _before = set()
            def _snap(h, _): _before.add(h); return True
            try: _w32g_oa.EnumWindows(_snap, None)
            except Exception: pass
            subprocess.Popen(_cmd, shell=True)
            _t2.sleep(1.5)
            # 找到新出現的視窗
            _last_opened_hwnd = 0
            try:
                _kw = app.lower().replace(".exe", "").split()[-1]
                _after = []
                def _new_win(h, _):
                    if h not in _before and _w32g_oa.IsWindowVisible(h) and _w32g_oa.GetWindowText(h).strip():
                        _after.append(h)
                    return True
                _w32g_oa.EnumWindows(_new_win, None)
                # 優先找包含 app 關鍵字的新視窗
                _matched = [h for h in _after if _kw in _w32g_oa.GetWindowText(h).lower()]
                _last_opened_hwnd = (_matched or _after or [0])[0]
                if _last_opened_hwnd:
                    import ctypes as _ct_oa
                    _ct_oa.windll.user32.SetForegroundWindow(_last_opened_hwnd)
                    _t2.sleep(0.3)
            except Exception:
                pass
            return {"ok": True, "message": f"已開啟並切換到：{app}，視窗已就緒，可以直接輸入文字", "screenshot": None}

        elif action == "scroll":
            scroll_amount = amount if direction == "up" else -amount
            if x is not None and y is not None:
                ax, ay = _resolve_coords(x, y, monitor)
                pyautogui.scroll(scroll_amount, x=ax, y=ay)
            else:
                pyautogui.scroll(scroll_amount)
            return {"ok": True, "message": f"已向{direction}滾動 {amount} 格", "screenshot": None}

        else:
            return {"ok": False, "message": f"未知動作：{action}", "screenshot": None}

    except Exception as e:
        return {"ok": False, "message": f"執行失敗：{str(e)}", "screenshot": None}


def execute_device_manager(action, name="", keyword=""):
    try:
        import subprocess
        if action == "list":
            q = f"| Where-Object {{$_.Name -like '*{keyword}*'}}" if keyword else ""
            result = subprocess.run(["powershell", "-Command",
                f"Get-PnpDevice {q} | Select-Object Status,Class,FriendlyName | Format-Table -AutoSize"],
                capture_output=True, text=True, encoding="utf-8", errors="replace")
            return f"🖥️ 裝置清單：\n{result.stdout.strip()[:2000]}"
        elif action == "enable":
            result = subprocess.run(["powershell", "-Command",
                f"Get-PnpDevice | Where-Object {{$_.FriendlyName -like '*{name}*'}} | Enable-PnpDevice -Confirm:$false"],
                capture_output=True, text=True)
            return f"✅ 已啟用裝置：{name}" if result.returncode == 0 else f"❌ 啟用失敗：{result.stderr}"
        elif action == "disable":
            result = subprocess.run(["powershell", "-Command",
                f"Get-PnpDevice | Where-Object {{$_.FriendlyName -like '*{name}*'}} | Disable-PnpDevice -Confirm:$false"],
                capture_output=True, text=True)
            return f"✅ 已停用裝置：{name}" if result.returncode == 0 else f"❌ 停用失敗：{result.stderr}"
    except Exception as e:
        return f"❌ 裝置管理員失敗：{e}"


def execute_dialog_auto(action, button_text="", window_title="", timeout=30):
    try:
        import win32gui, win32con, time as _t
        def _find_dialogs():
            dialogs = []
            def cb(hwnd, _):
                if win32gui.IsWindowVisible(hwnd):
                    title = win32gui.GetWindowText(hwnd); cls = win32gui.GetClassName(hwnd)
                    if cls in ("#32770", "TForm") or any(kw in title for kw in ["確認","警告","錯誤","提示","Dialog","Error","Warning","Confirm"]):
                        dialogs.append((hwnd, title, cls))
            win32gui.EnumWindows(cb, None); return dialogs
        def _click_button(hwnd, text):
            import win32api; result = []
            def cb(child, _):
                t = win32gui.GetWindowText(child)
                if text.lower() in t.lower() or not text:
                    if win32gui.GetClassName(child) == "Button":
                        win32gui.SetForegroundWindow(hwnd); win32api.SendMessage(child, win32con.BM_CLICK, 0, 0); result.append(t)
            win32gui.EnumChildWindows(hwnd, cb, None); return result
        if action == "list_dialogs":
            dialogs = _find_dialogs()
            if not dialogs: return "目前無對話框"
            return "\n".join(f"HWND:{h} 標題:{t} 類別:{c}" for h, t, c in dialogs)
        elif action in ("find_and_click", "wait_and_click"):
            start = _t.time(); btn = button_text or "確定"
            while True:
                dialogs = _find_dialogs()
                for hwnd, title, _ in dialogs:
                    if window_title and window_title.lower() not in title.lower(): continue
                    clicked = _click_button(hwnd, btn)
                    if clicked: return f"✅ 已點擊對話框「{title}」的「{clicked[0]}」按鈕"
                if action == "find_and_click": return "❌ 未找到符合的對話框"
                if _t.time() - start > timeout: return f"❌ 等待 {timeout} 秒仍未出現對話框"
                _t.sleep(1)
        return "未知動作"
    except Exception as e:
        return f"❌ dialog_auto 失敗：{e}"


def execute_disk_analyze(path="C:/", top=10):
    try:
        import psutil
        usage = psutil.disk_usage(path)
        result = f"磁碟：{path}\n總容量：{usage.total/1024**3:.1f} GB | 已使用：{usage.used/1024**3:.1f} GB ({usage.percent}%) | 可用：{usage.free/1024**3:.1f} GB\n\n"
        sizes = []
        for item in Path(path).iterdir():
            try:
                size = sum(f.stat().st_size for f in item.rglob("*") if f.is_file()) if item.is_dir() else item.stat().st_size
                sizes.append((size, str(item)))
            except Exception: pass
        sizes.sort(reverse=True)
        result += "\n".join(f"{s/1024**3:.2f} GB  {n}" for s, n in sizes[:top])
        return result
    except Exception as e:
        return f"❌ 磁碟分析失敗：{e}"


def execute_disk_backup(action, src="", dest=""):
    try:
        import tempfile, shutil
        tmp = Path(tempfile.gettempdir())
        if action == "list_temp":
            files = list(tmp.rglob("*"))
            total = sum(f.stat().st_size for f in files if f.is_file())
            return f"暫存資料夾：{tmp}\n檔案數：{len(files)}\n佔用空間：{total/1024/1024:.1f} MB"
        elif action == "clean_temp":
            count = 0
            for f in tmp.iterdir():
                try:
                    if f.is_file(): f.unlink(); count += 1
                    elif f.is_dir(): shutil.rmtree(f, ignore_errors=True); count += 1
                except Exception: pass
            return f"✅ 已清理 {count} 個暫存項目"
        elif action == "backup":
            out = Path(dest) / f"{Path(src).name}_{dt.dt.datetime.now().strftime('%Y%m%d_%H%M%S')}"
            shutil.make_archive(str(out), "zip", src)
            return f"✅ 備份完成：{out}.zip"
    except Exception as e:
        return f"❌ 磁碟操作失敗：{e}"


def execute_display(action, level=None):
    try:
        import screen_brightness_control as sbc
        if action == "brightness_get":
            b = sbc.get_brightness()
            return f"💡 亮度：{b}%"
        elif action == "brightness_set":
            sbc.set_brightness(max(0, min(100, int(level))))
            return f"✅ 亮度設定為 {level}%"
        elif action == "resolution":
            import subprocess
            result = subprocess.run(["powershell", "-Command",
                "Get-CimInstance Win32_VideoController | Select-Object CurrentHorizontalResolution,CurrentVerticalResolution,VideoModeDescription | Format-List"],
                capture_output=True, text=True)
            return f"🖥️ 解析度資訊：\n{result.stdout.strip()}"
    except Exception as e:
        return f"❌ 螢幕控制失敗：{e}"


def execute_doc_ai(action, path="", path2="", fields="", question="", url=""):
    try:
        import anthropic as _ant, base64, mimetypes
        from pathlib import Path as _P
        client = _ant.Anthropic(api_key=os.getenv("ANTHROPIC_API_KEY", ""))
        def _load_content(p, u=""):
            if u: return [{"type": "text", "text": f"請分析這個網址的內容：{u}"}]
            if not p or not _P(p).exists(): return [{"type": "text", "text": f"（檔案不存在：{p}）"}]
            ext = _P(p).suffix.lower()
            if ext in (".png", ".jpg", ".jpeg", ".gif", ".webp"):
                with open(p, "rb") as f: data = base64.b64encode(f.read()).decode()
                mime = mimetypes.guess_type(p)[0] or "image/jpeg"
                return [{"type": "image", "source": {"type": "base64", "media_type": mime, "data": data}}]
            elif ext == ".pdf":
                try:
                    import pdfplumber
                    with pdfplumber.open(p) as pdf: text = "\n".join(pg.extract_text() or "" for pg in pdf.pages[:10])
                except ImportError: text = f"（需安裝 pdfplumber 才能讀取 PDF：{p}）"
                return [{"type": "text", "text": text[:4000]}]
            else:
                with open(p, "r", encoding="utf-8", errors="ignore") as f: return [{"type": "text", "text": f.read()[:4000]}]
        content = _load_content(path, url)
        if action == "extract": prompt = f"請從以下內容提取這些欄位：{fields or '所有重要資訊'}。以 JSON 格式回傳。"
        elif action == "summarize": prompt = "請用繁體中文摘要這份文件的主要內容（200字以內）。"
        elif action == "classify": prompt = "請判斷這份文件的類型（如：發票、合約、報表、履歷等），並說明判斷依據。"
        elif action == "qa": prompt = f"根據文件內容回答：{question}"
        elif action == "compare":
            content2 = _load_content(path2); prompt = "請比較這兩份文件的差異，列出主要不同點。"
            content = content + [{"type": "text", "text": "---第二份文件---"}] + content2
        else: prompt = "請分析並說明這份文件的內容。"
        resp = client.messages.create(model="claude-sonnet-4-6", max_tokens=1024,
            messages=[{"role": "user", "content": content + [{"type": "text", "text": prompt}]}])
        return resp.content[0].text
    except Exception as e: return f"❌ doc_ai 失敗：{e}"


def execute_docker(action, name=""):
    try:
        import docker as _docker; client = _docker.from_env()
        if action == "list": return "\n".join(f"[{c.status}] {c.name} {c.image.tags}" for c in client.containers.list(all=True)) or "（無容器）"
        elif action == "start": client.containers.get(name).start(); return f"✅ {name} 已啟動"
        elif action == "stop": client.containers.get(name).stop(); return f"✅ {name} 已停止"
        elif action == "logs": return client.containers.get(name).logs(tail=50).decode(errors="replace")
        elif action == "images": return "\n".join(f"{img.tags} {img.short_id}" for img in client.images.list())
    except Exception as e: return f"❌ Docker 失敗：{e}"


def execute_document(action, path, content="", sheet=None):
    try:
        if action == "word_read":
            from docx import Document; doc = Document(path)
            return "\n".join(p.text for p in doc.paragraphs if p.text.strip()) or "（文件為空）"
        elif action == "word_write":
            from docx import Document
            doc = Document() if not Path(path).exists() else Document(path)
            doc.add_paragraph(content); doc.save(path)
            return f"已寫入：{path}"
        elif action == "excel_read":
            import openpyxl; wb = openpyxl.load_workbook(path)
            ws = wb[sheet] if sheet else wb.active
            rows = [[str(c.value or "") for c in row] for row in ws.iter_rows()]
            return "\n".join(["\t".join(r) for r in rows[:30]])
        elif action == "excel_write":
            import openpyxl, json
            wb = openpyxl.load_workbook(path) if Path(path).exists() else openpyxl.Workbook()
            ws = wb[sheet] if (sheet and sheet in wb.sheetnames) else wb.active
            data = json.loads(content)
            for row in data: ws.append(row)
            wb.save(path)
            return f"已寫入 {len(data)} 行到 {path}"
        elif action == "pdf_read":
            import fitz; doc = fitz.open(path)
            text = "\n".join(page.get_text() for page in doc)
            return text[:3000] if text else "（PDF 無可讀文字）"
        return "未知動作"
    except Exception as e:
        return f"文件操作失敗：{e}"


def execute_download_file(url, save_path=""):
    try:
        import requests
        if not save_path:
            fname = url.split("/")[-1].split("?")[0] or "download"
            save_path = str(Path("C:/Users/blue_/Desktop/測試檔案") / fname)
        r = requests.get(url, stream=True, timeout=60)
        r.raise_for_status()
        with open(save_path, "wb") as f:
            for chunk in r.iter_content(chunk_size=8192):
                f.write(chunk)
        return f"✅ 已下載：{save_path}"
    except Exception as e:
        return f"❌ 下載失敗：{e}"


def execute_dropbox(action, local, remote, token=""):
    try:
        import dropbox as dbx; tok = token or os.getenv("DROPBOX_TOKEN", "")
        if not tok: return "❌ 請設定 DROPBOX_TOKEN 環境變數"
        d = dbx.Dropbox(tok)
        if action == "upload":
            with open(local, "rb") as f: d.files_upload(f.read(), remote, mode=dbx.files.WriteMode.overwrite)
            return f"✅ 已上傳到 Dropbox：{remote}"
        elif action == "download": _, res = d.files_download(remote); Path(local).write_bytes(res.content); return f"✅ 已從 Dropbox 下載：{local}"
    except Exception as e: return f"❌ Dropbox 失敗：{e}"


def execute_email_read(host, user, password, folder="INBOX", count=5):
    try:
        import imapclient, email as _email
        from email.header import decode_header
        client = imapclient.IMAPClient(host, ssl=True)
        client.login(user, password); client.select_folder(folder)
        msgs = client.search(["ALL"])
        recent = msgs[-count:] if len(msgs) >= count else msgs
        results = []
        for uid in reversed(recent):
            raw = client.fetch([uid], ["RFC822"])[uid][b"RFC822"]
            msg = _email.message_from_bytes(raw)
            subj_raw, enc = decode_header(msg["Subject"])[0]
            subject = subj_raw.decode(enc or "utf-8") if isinstance(subj_raw, bytes) else subj_raw
            body = ""
            if msg.is_multipart():
                for part in msg.walk():
                    if part.get_content_type() == "text/plain":
                        body = part.get_payload(decode=True).decode(errors="replace")[:150]; break
            else:
                body = msg.get_payload(decode=True).decode(errors="replace")[:150]
            results.append(f"寄件人：{msg['From']}\n主旨：{subject}\n日期：{msg['Date']}\n{body}\n{'─'*30}")
        client.logout()
        return "\n".join(results) if results else "（收件匣為空）"
    except Exception as e:
        return f"❌ 讀取郵件失敗：{e}"


def execute_email_trigger(action, host="", user="", password="", filter_from="",
                          filter_subject="", duration=300, to="", subject="", body=""):
    try:
        if action == "send":
            import smtplib
            from email.mime.text import MIMEText
            smtp_host = host.replace("imap.", "smtp.") if host else "smtp.gmail.com"
            msg = MIMEText(body, "plain", "utf-8")
            msg["Subject"] = subject
            msg["From"] = user
            msg["To"] = to
            with smtplib.SMTP_SSL(smtp_host, 465) as s:
                s.login(user, password)
                s.send_message(msg)
            return f"✅ 郵件已發送至 {to}"

        import imaplib, email as _email, time as _t
        def _connect():
            m = imaplib.IMAP4_SSL(host or "imap.gmail.com")
            m.login(user, password)
            m.select("INBOX")
            return m

        def _fetch_recent(m, n=5):
            _, data = m.search(None, "ALL")
            ids = data[0].split()[-n:]
            mails = []
            for mid in reversed(ids):
                _, md = m.fetch(mid, "(RFC822)")
                msg = _email.message_from_bytes(md[0][1])
                sender = msg.get("From", "")
                subj = msg.get("Subject", "")
                if filter_from and filter_from.lower() not in sender.lower(): continue
                if filter_subject and filter_subject.lower() not in subj.lower(): continue
                body_text = ""
                if msg.is_multipart():
                    for part in msg.walk():
                        if part.get_content_type() == "text/plain":
                            body_text = part.get_payload(decode=True).decode("utf-8", errors="ignore")[:300]; break
                else:
                    body_text = msg.get_payload(decode=True).decode("utf-8", errors="ignore")[:300]
                mails.append(f"寄件人：{sender}\n主旨：{subj}\n內容：{body_text}")
            return mails

        if action == "check":
            m = _connect(); mails = _fetch_recent(m, 5); m.logout()
            return "\n---\n".join(mails) if mails else "📭 收件箱無符合郵件"
        elif action == "watch":
            m = _connect(); seen = set()
            _, data = m.search(None, "ALL")
            for mid in data[0].split(): seen.add(mid)
            results = []; start = _t.time()
            while _t.time() - start < duration:
                _t.sleep(10)
                try:
                    m.noop(); _, data = m.search(None, "ALL")
                    for mid in data[0].split():
                        if mid not in seen:
                            seen.add(mid); _, md = m.fetch(mid, "(RFC822)")
                            msg = _email.message_from_bytes(md[0][1])
                            results.append(f"📬 新郵件：{msg.get('From')} | {msg.get('Subject')}")
                except Exception: m = _connect()
            m.logout()
            return "\n".join(results) if results else f"監控 {duration} 秒內無新郵件"
        return "未知動作"
    except Exception as e:
        return f"❌ email_trigger 失敗：{e}"


def execute_emotion_detect(action, text="", image_path=""):
    """情緒偵測：從文字/臉部偵測情緒狀態"""

    if action == "from_text":
        try:
            pos_words = ["開心","高興","棒","好","讚","喜歡","愛","感謝","謝謝","哈哈","😊","😄","👍","❤️","爽","讚讚"]
            neg_words = ["難過","生氣","討厭","煩","累","痛","哭","傷心","憤怒","失望","😢","😡","😤","💔","爛","幹"]
            pos = sum(1 for w in pos_words if w in text)
            neg = sum(1 for w in neg_words if w in text)
            basic = "正面" if pos > neg else ("負面" if neg > pos else "中性")

            from anthropic import Anthropic
            c = Anthropic()
            resp = c.messages.create(
                model="claude-haiku-4-5-20251001", max_tokens=200,
                messages=[{"role": "user", "content": f"分析情緒（一行一項）：主要情緒、強度1-10、建議回應\n文字：{text}"}]
            )
            return f"情緒分析\n規則：{basic}（正面{pos}負面{neg}）\n\nAI分析：\n{resp.content[0].text}"
        except Exception as e:
            return f"分析失敗：{e}"

    elif action == "from_face":
        try:
            import base64, io as _io
            from PIL import Image as PILImg
            if image_path:
                img = PILImg.open(image_path)
            else:
                import pyautogui
                img = pyautogui.screenshot()
            buf = _io.BytesIO()
            img.save(buf, format="JPEG")
            b64 = base64.b64encode(buf.getvalue()).decode()

            from anthropic import Anthropic
            c = Anthropic()
            resp = c.messages.create(
                model="claude-sonnet-4-6", max_tokens=300,
                messages=[{"role": "user", "content": [
                    {"type": "image", "source": {"type": "base64", "media_type": "image/jpeg", "data": b64}},
                    {"type": "text", "text": "分析圖片中人臉情緒：臉數、每張臉的情緒（開心/悲傷/憤怒/驚訝/中性）、整體氛圍"}
                ]}]
            )
            return f"臉部情緒分析：\n{resp.content[0].text}"
        except Exception as e:
            return f"臉部分析失敗：{e}"

    return f"未知動作：{action}"


def execute_encrypt_file(action, path, password):
    try:
        from cryptography.hazmat.primitives.kdf.pbkdf2 import PBKDF2HMAC; from cryptography.hazmat.primitives import hashes; from cryptography.fernet import Fernet; import base64
        kdf = PBKDF2HMAC(algorithm=hashes.SHA256(), length=32, salt=b"xiaoniuma_salt_v1", iterations=100000)
        key = base64.urlsafe_b64encode(kdf.derive(password.encode())); f = Fernet(key); data = Path(path).read_bytes()
        if action == "encrypt": enc = f.encrypt(data); out = path + ".enc"; Path(out).write_bytes(enc); return f"✅ 已加密：{out}"
        elif action == "decrypt": dec = f.decrypt(data); out = path.replace(".enc", ".dec"); Path(out).write_bytes(dec); return f"✅ 已解密：{out}"
        return "未知動作"
    except Exception as e: return f"加密/解密失敗：{e}"


def execute_env_var(action, name="", value="", permanent=False):
    try:
        import os, subprocess
        if action == "get":
            v = os.environ.get(name, "")
            if not v:
                result = subprocess.run(["powershell", "-Command",
                    f"[System.Environment]::GetEnvironmentVariable('{name}','Machine')"],
                    capture_output=True, text=True)
                v = result.stdout.strip()
            return f"🌍 {name} = {v}" if v else f"⚠️ 環境變數 {name} 不存在"
        elif action == "set":
            os.environ[name] = value
            if permanent:
                subprocess.run(["powershell", "-Command",
                    f"[System.Environment]::SetEnvironmentVariable('{name}','{value}','User')"],
                    capture_output=True)
            return f"✅ 環境變數設定：{name}={value}{'（永久）' if permanent else '（本次）'}"
    except Exception as e:
        return f"❌ 環境變數失敗：{e}"


def execute_event_log(log="System", level="Error", count=10):
    try:
        import win32evtlog, win32evtlogutil, win32con
        level_map = {"Error": 1, "Warning": 2, "Information": 4, "All": 7}
        event_type = level_map.get(level, 1)
        hand = win32evtlog.OpenEventLog(None, log)
        flags = win32evtlog.EVENTLOG_BACKWARDS_READ | win32evtlog.EVENTLOG_SEQUENTIAL_READ
        events = []
        while len(events) < int(count):
            batch = win32evtlog.ReadEventLog(hand, flags, 0)
            if not batch:
                break
            for e in batch:
                if level == "All" or e.EventType == event_type:
                    try:
                        msg = win32evtlogutil.SafeFormatMessage(e, log)[:100]
                    except Exception:
                        msg = "(無法讀取訊息)"
                    events.append(f"[{e.TimeGenerated.Format()}] {e.SourceName}: {msg}")
                if len(events) >= int(count):
                    break
        win32evtlog.CloseEventLog(hand)
        if not events:
            return f"✅ {log} 中沒有 {level} 等級事件"
        return f"📋 {log} 事件記錄（{level}）：\n" + "\n".join(events)
    except Exception as e:
        return f"❌ 事件記錄讀取失敗：{e}"


def execute_excel_chart(path, sheet, chart_type="bar", title=""):
    try:
        import openpyxl
        from openpyxl.chart import BarChart, LineChart, PieChart, Reference
        wb = openpyxl.load_workbook(path)
        ws = wb[sheet] if sheet in wb.sheetnames else wb.active
        chart = {"bar": BarChart, "line": LineChart, "pie": PieChart}.get(chart_type, BarChart)()
        chart.title = title or sheet; chart.style = 10
        data = Reference(ws, min_col=2, min_row=1, max_row=ws.max_row, max_col=ws.max_column)
        chart.add_data(data, titles_from_data=True)
        chart.set_categories(Reference(ws, min_col=1, min_row=2, max_row=ws.max_row))
        ws.add_chart(chart, "A" + str(ws.max_row + 2))
        wb.save(path); return f"✅ 圖表已加入：{path}"
    except Exception as e:
        return f"❌ Excel 圖表失敗：{e}"


def execute_face_detect(image_path="", output=""):
    try:
        import cv2, numpy as np
        img = cv2.cvtColor(np.array(pyautogui.screenshot()), cv2.COLOR_RGB2BGR) if not image_path else cv2.imread(image_path)
        cascade = cv2.CascadeClassifier(cv2.data.haarcascades + "haarcascade_frontalface_default.xml")
        faces = cascade.detectMultiScale(cv2.cvtColor(img, cv2.COLOR_BGR2GRAY), 1.1, 4)
        for (x, y, w, h) in faces:
            cv2.rectangle(img, (x, y), (x+w, y+h), (0, 255, 0), 2)
        out_path = output or str(Path("C:/Users/blue_/Desktop/測試檔案") / f"faces_{dt.dt.datetime.now().strftime('%H%M%S')}.jpg")
        cv2.imwrite(out_path, img)
        return f"✅ 偵測到 {len(faces)} 張人臉：{out_path}"
    except Exception as e:
        return f"❌ 人臉偵測失敗：{e}"


def execute_face_recognize(action, name="", image_path="", output=""):
    try:
        import cv2 as _cv2, numpy as _np, json, os as _os
        from pathlib import Path as _P
        FACE_DB = str(_P.home() / ".face_db"); _os.makedirs(FACE_DB, exist_ok=True)
        try: import face_recognition as _fr
        except ImportError: return "❌ 請先安裝：pip install face-recognition opencv-python"
        def _capture_frame():
            cap = _cv2.VideoCapture(0); ret, frame = cap.read(); cap.release()
            if not ret: raise RuntimeError("無法開啟攝影機"); return frame
        if action == "capture":
            frame = _capture_frame(); out = output or str(_P("C:/Users/blue_/Desktop/測試檔案") / "face_capture.jpg")
            _cv2.imwrite(out, frame); return f"✅ 已拍照：{out}"
        elif action == "detect":
            img = _fr.load_image_file(image_path) if image_path else _capture_frame()
            if image_path: img = _cv2.cvtColor(_cv2.imread(image_path), _cv2.COLOR_BGR2RGB)
            locs = _fr.face_locations(img); return f"✅ 偵測到 {len(locs)} 個人臉"
        elif action == "enroll":
            if not name: return "❌ 需提供 name"
            img = _fr.load_image_file(image_path) if image_path else _capture_frame()
            if not image_path: img = _cv2.cvtColor(img, _cv2.COLOR_BGR2RGB)
            encs = _fr.face_encodings(img)
            if not encs: return "❌ 未偵測到人臉"
            _np.save(_os.path.join(FACE_DB, f"{name}.npy"), encs[0]); return f"✅ 已登記人臉：{name}"
        elif action == "recognize":
            known_encs, known_names = [], []
            for f in _P(FACE_DB).glob("*.npy"): known_encs.append(_np.load(str(f))); known_names.append(f.stem)
            if not known_encs: return "❌ 尚未登記任何人臉，請先用 enroll"
            img = _fr.load_image_file(image_path) if image_path else _capture_frame()
            if not image_path: img = _cv2.cvtColor(img, _cv2.COLOR_BGR2RGB)
            encs = _fr.face_encodings(img)
            if not encs: return "❌ 未偵測到人臉"
            results = []
            for enc in encs:
                matches = _fr.compare_faces(known_encs, enc); dists = _fr.face_distance(known_encs, enc)
                best = int(_np.argmin(dists))
                if matches[best]: results.append(f"✅ 識別為：{known_names[best]}（相似度 {(1-dists[best])*100:.0f}%）")
                else: results.append("❓ 未知人物")
            return "\n".join(results)
        return "未知動作"
    except Exception as e: return f"❌ face_recognize 失敗：{e}"


def execute_file_diff(file1, file2, output="", mode="unified"):
    try:
        import difflib
        text1 = Path(file1).read_text(encoding="utf-8", errors="replace").splitlines(keepends=True)
        text2 = Path(file2).read_text(encoding="utf-8", errors="replace").splitlines(keepends=True)
        if mode == "unified": diff = list(difflib.unified_diff(text1, text2, fromfile=file1, tofile=file2))
        elif mode == "context": diff = list(difflib.context_diff(text1, text2, fromfile=file1, tofile=file2))
        else: diff = [f"{'+ ' if a != b else '  '}{a}" for a, b in zip(text1, text2)]
        if not diff: return "✅ 兩個檔案內容完全相同"
        result = "".join(diff)
        if output:
            Path(output).write_text(result, encoding="utf-8")
            return f"✅ diff 已儲存：{output}\n（共 {len(diff)} 行差異）"
        return f"📄 檔案差異（{len(diff)} 行）：\n{result[:2000]}"
    except Exception as e:
        return f"❌ 檔案比較失敗：{e}"


def execute_file_system(action, path="", dest="", content="", keyword=""):
    try:
        if action == "list":
            p = Path(path or ".")
            items = sorted(p.iterdir())
            return "\n".join(("📁 " if i.is_dir() else "📄 ") + i.name for i in items)
        elif action == "read":
            return Path(path).read_text(encoding="utf-8", errors="replace")[:3000]
        elif action == "write":
            Path(path).write_text(content, encoding="utf-8")
            return f"已寫入：{path}"
        elif action == "delete":
            p = Path(path)
            shutil.rmtree(p) if p.is_dir() else p.unlink()
            return f"已刪除：{path}"
        elif action == "copy":
            shutil.copy2(path, dest)
            return f"已複製：{path} → {dest}"
        elif action == "move":
            shutil.move(path, dest)
            return f"已移動：{path} → {dest}"
        elif action == "search":
            results = list(Path(path).rglob(f"*{keyword}*"))
            return "\n".join(str(r) for r in results[:50]) or "找不到結果"
    except Exception as e:
        return f"執行失敗：{e}"


def execute_file_tools(action, path, dest="", pattern="", replacement="", ext=""):
    try:
        import re, filecmp, shutil
        if action == "batch_rename":
            files = [f for f in Path(path).iterdir() if f.is_file() and (not ext or f.suffix.lower()==ext.lower())]
            count = 0
            for f in sorted(files):
                new_name = re.sub(pattern, replacement, f.stem) + f.suffix
                if new_name != f.name: f.rename(f.parent / new_name); count += 1
            return f"✅ 已重新命名 {count} 個檔案"
        elif action == "sync":
            dest_path = Path(dest); dest_path.mkdir(parents=True, exist_ok=True)
            copied = 0
            for item in Path(path).rglob("*"):
                rel = item.relative_to(path); dst = dest_path / rel
                if item.is_dir(): dst.mkdir(parents=True, exist_ok=True)
                elif item.is_file() and (not dst.exists() or not filecmp.cmp(str(item), str(dst), shallow=False)):
                    shutil.copy2(str(item), str(dst)); copied += 1
            return f"✅ 同步完成：{copied} 個檔案更新"
    except Exception as e:
        return f"❌ 失敗：{e}"


def execute_file_transfer(action, source, dest=""):
    import zipfile
    if action == "zip":
        src = Path(source)
        with zipfile.ZipFile(dest, "w", zipfile.ZIP_DEFLATED) as zf:
            if src.is_dir():
                for f in src.rglob("*"):
                    zf.write(f, f.relative_to(src.parent))
            else:
                zf.write(src, src.name)
        return f"已壓縮：{source} → {dest}"
    elif action == "unzip":
        with zipfile.ZipFile(source, "r") as zf:
            zf.extractall(dest)
        return f"已解壓縮：{source} → {dest}"
    elif action == "download":
        if not dest:
            dest = str(Path("C:/Users/blue_/Desktop") / source.split("/")[-1].split("?")[0])
        r = requests.get(source, stream=True, timeout=60)
        r.raise_for_status()
        with open(dest, "wb") as f:
            for chunk in r.iter_content(chunk_size=8192):
                f.write(chunk)
        return f"已下載：{dest}"


def execute_file_trigger(folder, event, action, pattern="", target="", duration=60):
    try:
        import time as _t, fnmatch, shutil, subprocess as _sp
        from pathlib import Path as _P
        folder_path = _P(folder)
        if not folder_path.exists(): return f"❌ 資料夾不存在：{folder}"
        if action == "list":
            files = list(folder_path.iterdir())
            if pattern: files = [f for f in files if fnmatch.fnmatch(f.name, pattern)]
            return "\n".join(str(f) for f in files[:30]) or "資料夾為空"
        before = set(str(f) for f in folder_path.rglob("*"))
        triggered = []; start = _t.time()
        while _t.time() - start < duration:
            _t.sleep(2)
            after = set(str(f) for f in folder_path.rglob("*"))
            if event in ("created", "any"):
                new_files = after - before
                for f in new_files:
                    if pattern and not fnmatch.fnmatch(_P(f).name, pattern): continue
                    msg = f"[新增] {f}"
                    if action == "copy" and target: shutil.copy2(f, target); msg += f" → 已複製到 {target}"
                    elif action == "run" and target: _sp.Popen([target, f]); msg += f" → 已執行 {target}"
                    elif action == "notify": msg += " → 通知觸發"
                    triggered.append(msg)
            if event in ("deleted", "any"):
                removed = before - after
                for f in removed: triggered.append(f"[刪除] {f}")
            before = after
        return "\n".join(triggered) if triggered else f"監控 {duration} 秒內無 {event} 事件"
    except Exception as e:
        return f"❌ file_trigger 失敗：{e}"


def execute_file_watcher(action, name="", path="", events="all", command="", notify=True, _bot_send=None, _chat_id=None):
    global _file_watchers
    try:
        import threading
        if action == "list":
            if not _file_watchers: return "⚠️ 無執行中的監聽器"
            return "📁 檔案監聽器：\n" + "\n".join(f"- {k}: {v['path']}" for k,v in _file_watchers.items())
        elif action == "stop":
            if name in _file_watchers: _file_watchers[name]["observer"].stop(); del _file_watchers[name]; return f"✅ 已停止監聽：{name}"
            return f"⚠️ 找不到：{name}"
        elif action == "start":
            from watchdog.observers import Observer; from watchdog.events import FileSystemEventHandler
            class _Handler(FileSystemEventHandler):
                def _handle(self, event, etype):
                    if events != "all" and etype not in events: return
                    msg = f"📁 [{name}] {etype}：{event.src_path}"
                    if command: import subprocess; subprocess.Popen(command.replace("{path}", event.src_path), shell=True)
                    if notify and _bot_send and _chat_id:
                        import asyncio; asyncio.run_coroutine_threadsafe(_bot_send(_chat_id, msg), asyncio.get_event_loop())
                def on_created(self, e): self._handle(e, "created")
                def on_modified(self, e): self._handle(e, "modified")
                def on_deleted(self, e): self._handle(e, "deleted")
            observer = Observer(); observer.schedule(_Handler(), path, recursive=True); observer.start()
            _file_watchers[name] = {"path": path, "observer": observer}
            return f"✅ 已開始監聽：{name} → {path}（事件：{events}）"
    except Exception as e:
        return f"❌ 檔案監聽失敗：{e}"


def execute_find_image(template_path: str, confidence: float = 0.8) -> str:
    try:
        location = pyautogui.locateOnScreen(template_path, confidence=confidence)
        if location:
            cx, cy = pyautogui.center(location)
            return f"找到圖片，中心座標：({cx}, {cy})，區域：{location}"
        return "畫面上找不到該圖片"
    except Exception as e:
        return f"搜尋失敗：{str(e)}"


def execute_git(action, repo=".", message="", branch="master"):
    try:
        import git as _git; r = _git.Repo(repo)
        if action == "status": return r.git.status()
        elif action == "log": return "\n".join(f"{c.hexsha[:7]} [{c.authored_datetime.strftime('%m-%d %H:%M')}] {c.message.strip()[:60]}" for c in list(r.iter_commits())[:10])
        elif action == "pull": r.remotes.origin.pull(); return f"✅ Pull 完成"
        elif action == "add": r.git.add(A=True); return "✅ git add -A 完成"
        elif action == "commit": r.index.commit(message or "auto commit"); return f"✅ committed: {message}"
        elif action == "push": r.remotes.origin.push(branch); return f"✅ pushed to origin/{branch}"
        elif action == "diff": return r.git.diff()[:2000] or "（無變更）"
        return "未知動作"
    except Exception as e: return f"❌ Git 失敗：{e}"


def execute_global_hotkey(hotkey, command, duration=60.0):
    try:
        import keyboard as kb, time as t; triggered = []
        def on_trigger(): triggered.append(dt.dt.datetime.now().strftime("%H:%M:%S")); subprocess.run(command, shell=True)
        kb.add_hotkey(hotkey, on_trigger); t.sleep(duration); kb.remove_all_hotkeys()
        return f"✅ 快捷鍵 [{hotkey}] 共觸發 {len(triggered)} 次"
    except Exception as e: return f"❌ 快捷鍵監聽失敗：{e}"


def execute_goal_manager(action, goal="", goal_id="", steps="", priority="normal"):
    """目標管理系統"""
    GOALS_DB = Path(__file__).parent / "goals.db"
    def init_db():
        conn = sqlite3.connect(GOALS_DB)
        conn.execute("""CREATE TABLE IF NOT EXISTS goals (id INTEGER PRIMARY KEY AUTOINCREMENT, title TEXT NOT NULL, steps TEXT, status TEXT DEFAULT 'pending', priority TEXT DEFAULT 'normal', progress INTEGER DEFAULT 0, created_at TIMESTAMP DEFAULT CURRENT_TIMESTAMP, updated_at TIMESTAMP DEFAULT CURRENT_TIMESTAMP)""")
        conn.commit(); conn.close()
    init_db()
    if action == "add":
        try:
            if not steps:
                from anthropic import Anthropic; c = Anthropic()
                resp = c.messages.create(model="claude-sonnet-4-6", max_tokens=800, messages=[{"role": "user", "content": f"將以下目標分解為5-8個具體執行步驟，每步驟一行：\n\n目標：{goal}"}])
                steps = resp.content[0].text
            conn = sqlite3.connect(GOALS_DB); cur = conn.execute("INSERT INTO goals (title, steps, priority) VALUES (?, ?, ?)", (goal, steps, priority))
            gid = cur.lastrowid; conn.commit(); conn.close()
            return f"✅ 目標建立 (ID:{gid})\n{goal}\n\n步驟：\n{steps}"
        except Exception as e: return f"建立失敗：{e}"
    elif action == "list":
        conn = sqlite3.connect(GOALS_DB); rows = conn.execute("SELECT id, title, status, priority, progress FROM goals ORDER BY id DESC").fetchall(); conn.close()
        if not rows: return "目前沒有任何目標"
        icons = {"pending": "⏳", "in_progress": "🔄", "completed": "✅", "failed": "❌"}
        return "📋 目標清單：\n" + "\n".join(f"  {icons.get(r[2],'❓')} [{r[0]}] {r[1]} ({r[3]}) {r[4]}%" for r in rows)
    elif action == "detail":
        conn = sqlite3.connect(GOALS_DB); row = conn.execute("SELECT * FROM goals WHERE id=?", (goal_id,)).fetchone(); conn.close()
        if not row: return f"找不到目標 {goal_id}"
        return f"目標[{row[0]}]：{row[1]}\n狀態：{row[3]} | 優先：{row[4]} | 進度：{row[5]}%\n\n步驟：\n{row[2]}"
    elif action == "update_status":
        conn = sqlite3.connect(GOALS_DB); conn.execute("UPDATE goals SET status=?, updated_at=CURRENT_TIMESTAMP WHERE id=?", (steps, goal_id)); conn.commit(); conn.close()
        return f"✅ 目標 {goal_id} 狀態 → {steps}"
    elif action == "set_progress":
        conn = sqlite3.connect(GOALS_DB); conn.execute("UPDATE goals SET progress=?, updated_at=CURRENT_TIMESTAMP WHERE id=?", (int(steps or 0), goal_id)); conn.commit(); conn.close()
        return f"✅ 目標 {goal_id} 進度 → {steps}%"
    elif action == "delete":
        conn = sqlite3.connect(GOALS_DB); conn.execute("DELETE FROM goals WHERE id=?", (goal_id,)); conn.commit(); conn.close()
        return f"✅ 目標 {goal_id} 已刪除"
    elif action == "next_task":
        conn = sqlite3.connect(GOALS_DB); row = conn.execute("SELECT * FROM goals WHERE status IN ('pending','in_progress') ORDER BY priority DESC, id ASC LIMIT 1").fetchone(); conn.close()
        if not row: return "所有目標已完成或無目標"
        return f"🤖 下一個待執行目標：\n[{row[0]}] {row[1]}\n\n步驟：\n{row[2][:500] if row[2] else '待規劃'}"
    return f"未知動作：{action}"


def execute_hardware():
    try:
        import psutil
        battery = psutil.sensors_battery()
        bat = f"{battery.percent:.0f}% {'充電中' if battery.power_plugged else '使用電池'}" if battery else "無電池"
        gpu_str = ""
        try:
            import GPUtil
            gpus = GPUtil.getGPUs()
            gpu_str = "\n".join(f"GPU {g.name}: 使用率 {g.load*100:.0f}% | 記憶體 {g.memoryUsed:.0f}/{g.memoryTotal:.0f}MB | 溫度 {g.temperature}°C" for g in gpus)
        except Exception:
            gpu_str = "GPU：未偵測到 NVIDIA GPU"
        return f"🔋 電池：{bat}\n{gpu_str}"
    except Exception as e:
        return f"❌ 硬體監控失敗：{e}"


def execute_http_server(action, port=9876, password=""):
    global _http_control_proc
    try:
        import socket, subprocess as _sp, sys, os as _os
        if action == "start":
            if _http_control_proc and _http_control_proc.poll() is None:
                ip = socket.gethostbyname(socket.gethostname())
                return f"✅ HTTP 控制伺服器已在運行\n網址：http://{ip}:{port}/"
            script = f"""
import http.server, urllib.parse, json, subprocess, sys, os
PASSWORD = "{password}"
class H(http.server.BaseHTTPRequestHandler):
    def do_GET(self):
        if PASSWORD and self.headers.get('X-Password','') != PASSWORD:
            p = urllib.parse.urlparse(self.path)
            q = urllib.parse.parse_qs(p.query)
            if q.get('pw',[''])[0] != PASSWORD:
                self.send_response(403); self.end_headers(); self.wfile.write(b'Unauthorized'); return
        self.send_response(200); self.send_header('Content-Type','text/html; charset=utf-8'); self.end_headers()
        html = '''<html><body><h2>小牛馬遠端控制</h2><form method=post><input name=cmd size=60 placeholder="輸入指令"><button>執行</button></form></body></html>'''
        self.wfile.write(html.encode())
    def do_POST(self):
        l = int(self.headers.get('Content-Length',0))
        body = urllib.parse.parse_qs(self.rfile.read(l).decode())
        cmd = body.get('cmd',[''])[0]
        self.send_response(200); self.send_header('Content-Type','text/plain; charset=utf-8'); self.end_headers()
        if cmd:
            try: out = subprocess.check_output(cmd, shell=True, text=True, stderr=subprocess.STDOUT, timeout=10)
            except Exception as e: out = str(e)
            self.wfile.write(out.encode('utf-8','ignore'))
    def log_message(self, *a): pass
http.server.HTTPServer(('0.0.0.0', {port}), H).serve_forever()
"""
            _http_control_proc = _sp.Popen([sys.executable, "-c", script], creationflags=0x00000008)
            ip = socket.gethostbyname(socket.gethostname())
            return f"✅ HTTP 控制伺服器已啟動\n網址：http://{ip}:{port}/\n（手機在同一 WiFi 下可存取）"
        elif action == "stop":
            if _http_control_proc: _http_control_proc.terminate(); _http_control_proc = None
            return "✅ 已停止"
        elif action == "status":
            running = _http_control_proc and _http_control_proc.poll() is None
            return f"狀態：{'運行中' if running else '已停止'}"
        elif action == "get_url":
            ip = socket.gethostbyname(socket.gethostname())
            return f"http://{ip}:{port}/"
        return "未知動作"
    except Exception as e:
        return f"❌ http_server 失敗：{e}"


def execute_image_edit(action, path, *params):
    try:
        from PIL import Image, ImageDraw, ImageFont
        img = Image.open(path)
        if action == "crop":
            x, y, w, h = int(params[0]), int(params[1]), int(params[2]), int(params[3])
            img = img.crop((x, y, x+w, y+h))
        elif action == "resize":
            w, h = int(params[0]), int(params[1])
            img = img.resize((w, h))
        elif action == "text":
            x, y = int(params[0]), int(params[1])
            text = " ".join(params[2:]) if len(params) > 2 else ""
            draw = ImageDraw.Draw(img)
            draw.text((x, y), text, fill="red")
        elif action == "merge":
            img2 = Image.open(params[0])
            merged = Image.new("RGB", (img.width + img2.width, max(img.height, img2.height)))
            merged.paste(img, (0, 0))
            merged.paste(img2, (img.width, 0))
            img = merged
        out = path.replace(".", f"_{action}.")
        img.save(out)
        return f"✅ 圖片已儲存：{out}"
    except Exception as e:
        return f"圖片編輯失敗：{e}"


def execute_image_tools(action, path="", quality=75, width=0, height=0, target_lang="zh-TW"):
    try:
        if action == "compress":
            from PIL import Image
            img = Image.open(path)
            if img.mode in ("RGBA","P"): img = img.convert("RGB")
            out = path.replace(".", f"_q{quality}.")
            img.save(out, optimize=True, quality=quality)
            orig = Path(path).stat().st_size; new = Path(out).stat().st_size
            return f"✅ {orig//1024}KB → {new//1024}KB（節省 {(1-new/orig)*100:.1f}%）：{out}"
        elif action == "batch":
            from PIL import Image as _Img
            folder = Path(path); out_dir = folder / "output"; out_dir.mkdir(exist_ok=True)
            count = 0
            for f in list(folder.glob("*.jpg")) + list(folder.glob("*.png")) + list(folder.glob("*.jpeg")):
                img = _Img.open(f)
                if width and height: img = img.resize((width, height))
                if img.mode in ("RGBA","P"): img = img.convert("RGB")
                img.save(str(out_dir/f.name), optimize=True, quality=quality); count += 1
            return f"✅ 批次處理 {count} 張 → {out_dir}"
        elif action == "ocr_translate":
            import easyocr, numpy as np
            from PIL import Image
            from deep_translator import GoogleTranslator
            reader = easyocr.Reader(["ch_tra","en"], gpu=False)
            img = Image.open(path) if path else pyautogui.screenshot()
            text = " ".join(r[1] for r in reader.readtext(np.array(img)))
            if not text.strip(): return "❌ 未辨識到文字"
            translated = GoogleTranslator(source="auto", target=target_lang).translate(text)
            return f"原文：{text[:200]}\n翻譯：{translated}"
    except Exception as e:
        return f"❌ 失敗：{e}"


def execute_ime_switch(action):
    try:
        import ctypes, win32gui, win32api, win32con
        hwnd = win32gui.GetForegroundWindow()
        if action == "status":
            lid = ctypes.windll.user32.GetKeyboardLayout(0) & 0xFFFF
            return f"目前輸入法 LCID：{hex(lid)} ({'中文' if lid in (0x0804,0x0404,0x0C04) else '英文' if lid == 0x0409 else '其他'})"
        elif action == "switch_en":
            hkl = win32api.LoadKeyboardLayout("00000409", 1); win32api.PostMessage(hwnd, win32con.WM_INPUTLANGCHANGEREQUEST, 0, hkl)
            return "✅ 已切換至英文輸入"
        elif action == "switch_zh":
            hkl = win32api.LoadKeyboardLayout("00000804", 1); win32api.PostMessage(hwnd, win32con.WM_INPUTLANGCHANGEREQUEST, 0, hkl)
            return "✅ 已切換至中文輸入"
        elif action == "toggle": import pyautogui; pyautogui.hotkey("shift", "alt"); return "✅ 已切換輸入法"
        return "未知動作"
    except Exception as e:
        return f"❌ ime_switch 失敗：{e}"


def execute_knowledge_base(action, content="", query="", tag="", kb_id=""):
    """知識庫"""
    KB_DB = Path(__file__).parent / "knowledge_base.db"
    def init_kb():
        conn = sqlite3.connect(KB_DB)
        conn.execute("""CREATE TABLE IF NOT EXISTS knowledge (id INTEGER PRIMARY KEY AUTOINCREMENT, content TEXT NOT NULL, summary TEXT, tags TEXT DEFAULT '', created_at TIMESTAMP DEFAULT CURRENT_TIMESTAMP)""")
        conn.commit(); conn.close()
    init_kb()
    if action == "add":
        try:
            summary = content[:150] if len(content) > 150 else content
            conn = sqlite3.connect(KB_DB); cur = conn.execute("INSERT INTO knowledge (content, summary, tags) VALUES (?, ?, ?)", (content, summary, tag))
            kid = cur.lastrowid; conn.commit(); conn.close()
            return f"✅ 知識儲存 (ID:{kid})\n標籤：{tag}\n摘要：{summary[:100]}"
        except Exception as e: return f"儲存失敗：{e}"
    elif action == "search":
        try:
            conn = sqlite3.connect(KB_DB)
            rows = conn.execute("SELECT id, summary, tags, created_at FROM knowledge WHERE content LIKE ? OR tags LIKE ? ORDER BY id DESC LIMIT 10", (f"%{query}%", f"%{query}%")).fetchall()
            conn.close()
            if not rows: return f"找不到「{query}」相關知識"
            return f"🔍 找到 {len(rows)} 筆：\n" + "\n".join(f"  [{r[0]}] {r[1][:100]} [{r[2]}]" for r in rows)
        except Exception as e: return f"搜尋失敗：{e}"
    elif action == "get":
        conn = sqlite3.connect(KB_DB); row = conn.execute("SELECT * FROM knowledge WHERE id=?", (kb_id,)).fetchone(); conn.close()
        if not row: return f"找不到知識 {kb_id}"
        return f"[{row[0]}] {row[3] or '無標籤'} | {row[4]}\n\n{row[1]}"
    elif action == "list":
        conn = sqlite3.connect(KB_DB)
        rows = conn.execute("SELECT id, summary, tags, created_at FROM knowledge ORDER BY id DESC LIMIT 20").fetchall()
        total = conn.execute("SELECT COUNT(*) FROM knowledge").fetchone()[0]; conn.close()
        if not rows: return "知識庫空白"
        return f"📚 知識庫（共{total}條）：\n" + "\n".join(f"  [{r[0]}] {r[1][:80]} [{r[2]}]" for r in rows)
    elif action == "delete":
        conn = sqlite3.connect(KB_DB); conn.execute("DELETE FROM knowledge WHERE id=?", (kb_id,)); conn.commit(); conn.close()
        return f"✅ 知識 {kb_id} 已刪除"
    elif action == "stats":
        conn = sqlite3.connect(KB_DB)
        total = conn.execute("SELECT COUNT(*) FROM knowledge").fetchone()[0]
        tags = conn.execute("SELECT tags, COUNT(*) FROM knowledge GROUP BY tags ORDER BY COUNT(*) DESC LIMIT 10").fetchall(); conn.close()
        return f"📊 知識庫統計\n總條目：{total}\n\n標籤分佈：\n" + "\n".join(f"  {t[0] or '無標籤'}: {t[1]}條" for t in tags)
    return f"未知動作：{action}"


def execute_lan_scan(action, subnet="", host="", port=80):
    try:
        import socket, subprocess as _sp, concurrent.futures
        if action == "get_local_ip":
            s = socket.socket(socket.AF_INET, socket.SOCK_DGRAM); s.connect(("8.8.8.8", 80))
            ip = s.getsockname()[0]; s.close()
            return f"本機 IP：{ip}"
        elif action == "ping_sweep":
            if not subnet:
                s = socket.socket(socket.AF_INET, socket.SOCK_DGRAM); s.connect(("8.8.8.8", 80))
                local_ip = s.getsockname()[0]; s.close()
                subnet = ".".join(local_ip.split(".")[:3]) + ".0/24"
            base = ".".join(subnet.split(".")[:3])
            def ping(i):
                ip = f"{base}.{i}"
                r = _sp.run(["ping", "-n", "1", "-w", "300", ip], capture_output=True, text=True)
                if "TTL=" in r.stdout or "ttl=" in r.stdout:
                    try: hostname = socket.gethostbyaddr(ip)[0]
                    except Exception: hostname = "unknown"
                    return f"{ip} ({hostname})"
                return None
            with concurrent.futures.ThreadPoolExecutor(max_workers=50) as ex:
                results = ex.map(ping, range(1, 255))
            online = [r for r in results if r]
            return f"區域網路掃描結果（{subnet}）：\n" + "\n".join(online) if online else "無裝置回應"
        elif action == "port_check":
            s = socket.socket(socket.AF_INET, socket.SOCK_STREAM); s.settimeout(2)
            result = s.connect_ex((host, int(port))); s.close()
            return f"{'✅' if result == 0 else '❌'} {host}:{port} {'開啟' if result == 0 else '關閉'}"
        elif action == "scan":
            s = socket.socket(socket.AF_INET, socket.SOCK_DGRAM); s.connect(("8.8.8.8", 80))
            local_ip = s.getsockname()[0]; s.close()
            base = ".".join(local_ip.split(".")[:3])
            common_ports = [21, 22, 80, 443, 3389, 8080, 1883]
            def check(i):
                ip = f"{base}.{i}"
                r = _sp.run(["ping", "-n", "1", "-w", "200", ip], capture_output=True)
                if b"TTL=" in r.stdout or b"ttl=" in r.stdout:
                    open_ports = []
                    for p in common_ports:
                        sock = socket.socket(); sock.settimeout(0.3)
                        if sock.connect_ex((ip, p)) == 0: open_ports.append(str(p))
                        sock.close()
                    try: hn = socket.gethostbyaddr(ip)[0]
                    except Exception: hn = "?"
                    return f"{ip} ({hn}) 開放 port: {','.join(open_ports) or '無'}"
                return None
            with concurrent.futures.ThreadPoolExecutor(max_workers=50) as ex:
                res = list(ex.map(check, range(1, 255)))
            found = [r for r in res if r]
            return "\n".join(found) if found else "未找到裝置"
        return "未知動作"
    except Exception as e:
        return f"❌ lan_scan 失敗：{e}"


def execute_lookup(action, ip="", amount=1.0, from_cur="USD", to_cur="TWD"):
    try:
        if action == "ip":
            d = requests.get(f"http://ip-api.com/json/{ip}?lang=zh-TW", timeout=10).json()
            if d.get("status") == "fail": return f"❌ {d.get('message')}"
            return f"IP：{d['query']}\n國家：{d['country']}\n城市：{d['city']}\nISP：{d['isp']}\n時區：{d['timezone']}\n座標：{d['lat']},{d['lon']}"
        elif action == "currency":
            d = requests.get(f"https://api.frankfurter.app/latest?amount={amount}&from={from_cur.upper()}&to={to_cur.upper()}", timeout=10).json()
            rate = d["rates"].get(to_cur.upper())
            return f"💱 {amount} {from_cur.upper()} = {rate:.4f} {to_cur.upper()}（{d.get('date')}）"
    except Exception as e: return f"❌ 查詢失敗：{e}"


def execute_macro(action, name="", repeat=1, duration=10.0):
    try:
        import keyboard, json, time
        global _macro_store
        macro_file = Path.home() / ".claude_macros.json"
        if macro_file.exists(): _macro_store = json.loads(macro_file.read_text())
        if action == "record_start":
            keyboard.start_recording(); time.sleep(float(duration)); recorded = keyboard.stop_recording()
            _macro_store[name] = [{"type": e.event_type, "name": e.name, "time": e.time} for e in recorded]
            macro_file.write_text(json.dumps(_macro_store))
            return f"✅ 巨集 '{name}' 錄製完成（{len(recorded)} 個事件，{duration}s）"
        elif action == "play":
            if name not in _macro_store: return f"⚠️ 找不到巨集：{name}"
            events = _macro_store[name]
            for _ in range(int(repeat)):
                prev_time = events[0]["time"] if events else 0
                for e in events:
                    delay = max(0, e["time"] - prev_time); time.sleep(min(delay, 0.5)); prev_time = e["time"]
                    if e["type"] == "down": keyboard.press(e["name"])
                    elif e["type"] == "up": keyboard.release(e["name"])
            return f"✅ 巨集 '{name}' 已回放 {repeat} 次"
        elif action == "list":
            if not _macro_store: return "⚠️ 無已儲存巨集"
            return "📋 已儲存巨集：\n" + "\n".join(f"- {k}（{len(v)} 事件）" for k,v in _macro_store.items())
        elif action == "delete":
            if name in _macro_store: del _macro_store[name]; macro_file.write_text(json.dumps(_macro_store)); return f"✅ 已刪除巨集：{name}"
            return f"⚠️ 找不到巨集：{name}"
    except Exception as e: return f"❌ 巨集操作失敗：{e}"


def execute_manage_schedule(action: str, name: str = "", time: str = "", script: str = "") -> str:
    try:
        if action == "list":
            result = subprocess.run(
                [SCHTASKS, "/Query", "/FO", "CSV", "/NH"],
                capture_output=True, text=True, encoding="cp950", errors="replace"
            )
            lines = [l for l in result.stdout.strip().splitlines() if l.strip()]
            out = "目前排程任務：\n"
            for line in lines:
                parts = line.strip('"').split('","')
                if len(parts) >= 3:
                    t_name = parts[0].replace("\\", "").strip()
                    next_run = parts[1].strip()
                    status = parts[2].strip()
                    out += f"• {t_name} | 下次執行：{next_run} | {status}\n"
            return out.strip()

        elif action == "add":
            subprocess.run([SCHTASKS, "/Create", "/TN", name,
                            "/TR", f"pythonw {script}",
                            "/SC", "DAILY", "/ST", time, "/F"],
                           capture_output=True)
            ps = (
                f"$t = Get-ScheduledTask -TaskName '{name}';"
                f"$t.Settings.WakeToRun = $true;"
                f"$t.Settings.DisallowStartIfOnBatteries = $false;"
                f"$t.Settings.StopIfGoingOnBatteries = $false;"
                f"Set-ScheduledTask -TaskName '{name}' -Settings $t.Settings | Out-Null"
            )
            subprocess.run(["powershell.exe", "-Command", ps], capture_output=True)
            return f"排程 [{name}] 已建立，每天 {time} 執行，電腦待機也會自動喚醒。"

        elif action == "delete":
            result = subprocess.run([SCHTASKS, "/Delete", "/TN", name, "/F"],
                                    capture_output=True, text=True, encoding="cp950", errors="replace")
            if result.returncode == 0:
                return f"排程 [{name}] 已刪除。"
            else:
                return f"刪除失敗：{result.stderr.strip()}"

        elif action == "bot_status":
            result = subprocess.run(
                ["powershell.exe", "-Command",
                 "Get-Process pythonw -ErrorAction SilentlyContinue | Measure-Object | Select-Object -ExpandProperty Count"],
                capture_output=True, text=True
            )
            count = int(result.stdout.strip()) if result.stdout.strip().isdigit() else 0
            return f"Bot 執行中（{count} 個程序）✅" if count > 0 else "Bot 未執行 ❌"

        elif action == "bot_restart":
            subprocess.run(["powershell.exe", "-Command",
                            "Stop-Process -Name pythonw -Force -ErrorAction SilentlyContinue"])
            import time as t
            t.sleep(1)
            subprocess.Popen(["pythonw", BOT_SCRIPT], cwd=str(Path(BOT_SCRIPT).parent))
            return "Bot 已重啟 ✅"

        else:
            return f"未知動作：{action}"

    except Exception as e:
        return f"執行失敗：{str(e)}"


def execute_media(action, device_name=""):
    try:
        import keyboard
        key_map = {
            "play_pause": "play/pause media",
            "next": "next track",
            "prev": "previous track",
            "stop": "stop media",
            "volume_up": "volume up",
            "volume_down": "volume down",
            "mute": "volume mute"
        }
        if action in key_map:
            keyboard.send(key_map[action])
            return f"✅ 媒體控制：{action}"
        elif action == "list_devices":
            from pycaw.pycaw import AudioUtilities
            devices = AudioUtilities.GetAllDevices()
            names = [d.FriendlyName for d in devices if d.FriendlyName]
            return "🔊 音訊裝置：\n" + "\n".join(f"- {n}" for n in names)
        elif action == "switch_device":
            import subprocess
            result = subprocess.run(["powershell", "-Command",
                f"Get-AudioDevice -List | Where-Object {{$_.Name -like '*{device_name}*'}} | Set-AudioDevice"],
                capture_output=True, text=True)
            if result.returncode == 0:
                return f"✅ 已切換至：{device_name}"
            return f"⚠️ 切換失敗（需安裝 AudioDeviceCmdlets）：{result.stderr}"
    except Exception as e:
        return f"❌ 媒體控制失敗：{e}"


def execute_monitor_config():
    try:
        from screeninfo import get_monitors
        return "\n".join(f"{'主螢幕' if m.is_primary else '副螢幕'} {m.width}x{m.height} @({m.x},{m.y}) {m.name}" for m in get_monitors())
    except Exception as e: return f"❌ 取得螢幕資訊失敗：{e}"


def execute_mouse_record(action, name="", duration=10.0, repeat=1, speed=1.0):
    try:
        import json, time
        from pathlib import Path as _Path
        record_file = _Path.home() / ".claude_mouse_macros.json"
        store = json.loads(record_file.read_text()) if record_file.exists() else {}
        if action == "list":
            if not store: return "⚠️ 無已儲存的滑鼠巨集"
            return "🖱️ 滑鼠巨集：\n" + "\n".join(f"- {k}（{len(v)} 個事件）" for k,v in store.items())
        elif action == "delete":
            if name in store: del store[name]; record_file.write_text(json.dumps(store)); return f"✅ 已刪除：{name}"
            return f"⚠️ 找不到：{name}"
        elif action == "start":
            from pynput import mouse, keyboard as kb
            events = []; start_time = time.time()
            def on_move(x, y): events.append({"t": time.time()-start_time, "type":"move", "x":x, "y":y})
            def on_click(x, y, btn, pressed): events.append({"t": time.time()-start_time, "type":"click", "x":x, "y":y, "btn":str(btn), "pressed":pressed})
            def on_scroll(x, y, dx, dy): events.append({"t": time.time()-start_time, "type":"scroll", "x":x, "y":y, "dx":dx, "dy":dy})
            m_listener = mouse.Listener(on_move=on_move, on_click=on_click, on_scroll=on_scroll)
            m_listener.start(); time.sleep(float(duration)); m_listener.stop()
            store[name] = events; record_file.write_text(json.dumps(store))
            return f"✅ 滑鼠巨集 '{name}' 錄製完成（{len(events)} 個事件，{duration}s）"
        elif action == "play":
            if name not in store: return f"⚠️ 找不到：{name}"
            from pynput import mouse as m; controller = m.Controller(); events = store[name]; spd = float(speed)
            for _ in range(int(repeat)):
                prev_t = 0.0
                for e in events:
                    delay = (e["t"] - prev_t) / spd
                    if delay > 0: time.sleep(min(delay, 0.5)); prev_t = e["t"]
                    if e["type"] == "move": controller.position = (e["x"], e["y"])
                    elif e["type"] == "click":
                        btn = m.Button.left if "left" in e["btn"] else m.Button.right
                        if e["pressed"]: controller.press(btn)
                        else: controller.release(btn)
                    elif e["type"] == "scroll": controller.scroll(e["dx"], e["dy"])
            return f"✅ 滑鼠巨集 '{name}' 回放 {repeat} 次完成"
    except Exception as e:
        return f"❌ 滑鼠巨集失敗：{e}"


def execute_mqtt(action, broker, port=1883, topic="", message="", duration=10, username="", password=""):
    try:
        import paho.mqtt.client as _mqtt, time as _t
    except ImportError: return "❌ 請先安裝：pip install paho-mqtt"
    try:
        received = []; connected = [False]
        def on_connect(c, ud, flags, rc): connected[0] = rc == 0
        def on_message(c, ud, msg): received.append(f"[{msg.topic}] {msg.payload.decode('utf-8','ignore')}")
        client = _mqtt.Client(); client.on_connect = on_connect; client.on_message = on_message
        if username: client.username_pw_set(username, password)
        client.connect(broker, int(port), 60); client.loop_start(); _t.sleep(1)
        if not connected[0]: return f"❌ 無法連線至 {broker}:{port}"
        if action == "test_connect": client.disconnect(); return f"✅ 成功連線 {broker}:{port}"
        elif action == "publish": client.publish(topic, message); _t.sleep(0.5); client.disconnect(); return f"✅ 已發布到 {topic}：{message}"
        elif action == "subscribe": client.subscribe(topic); _t.sleep(duration); client.disconnect(); return "\n".join(received) if received else f"訂閱 {topic} 共 {duration} 秒，無訊息"
        return "未知動作"
    except Exception as e: return f"❌ mqtt 失敗：{e}"


def execute_multi_deploy(action, remote_host="", remote_user="", remote_pass="", remote_path="/tmp/niu_bot"):
    """多機器部署：將bot部署到遠端伺服器同步運行"""
    try: import paramiko
    except ImportError: return "需要安裝 paramiko: pip install paramiko"
    def get_ssh():
        ssh = paramiko.SSHClient(); ssh.set_missing_host_key_policy(paramiko.AutoAddPolicy())
        ssh.connect(remote_host, username=remote_user, password=remote_pass, timeout=15); return ssh
    if action == "deploy":
        try:
            ssh = get_ssh(); sftp = ssh.open_sftp()
            bot_path = Path(__file__); env_path = bot_path.parent / ".env"
            ssh.exec_command(f"mkdir -p {remote_path}"); time.sleep(0.5)
            sftp.put(str(bot_path), f"{remote_path}/bot.py")
            if env_path.exists(): sftp.put(str(env_path), f"{remote_path}/.env")
            sftp.close()
            for cmd in ["pip install python-telegram-bot anthropic requests python-dotenv feedparser -q",
                f"pkill -f '{remote_path}/bot.py' 2>/dev/null; true",
                f"nohup python {remote_path}/bot.py > {remote_path}/bot.log 2>&1 &"]:
                ssh.exec_command(cmd); time.sleep(1)
            ssh.close()
            return f"✅ Bot 已部署到 {remote_host}:{remote_path}"
        except Exception as e: return f"部署失敗：{e}"
    elif action == "status":
        try:
            ssh = get_ssh()
            _, stdout, _ = ssh.exec_command(f"pgrep -f '{remote_path}/bot.py' && echo RUNNING || echo STOPPED")
            status = stdout.read().decode().strip(); ssh.close()
            return f"{remote_host} Bot：{status}"
        except Exception as e: return f"查詢失敗：{e}"
    elif action == "sync":
        try:
            ssh = get_ssh(); sftp = ssh.open_sftp()
            sftp.put(str(Path(__file__)), f"{remote_path}/bot.py"); sftp.close()
            ssh.exec_command(f"pkill -f '{remote_path}/bot.py'; sleep 1; nohup python {remote_path}/bot.py > {remote_path}/bot.log 2>&1 &")
            ssh.close()
            return f"✅ 技能已同步到 {remote_host} 並重啟"
        except Exception as e: return f"同步失敗：{e}"
    elif action == "log":
        try:
            ssh = get_ssh()
            _, stdout, _ = ssh.exec_command(f"tail -50 {remote_path}/bot.log")
            log = stdout.read().decode(); ssh.close()
            return f"遠端日誌 ({remote_host}):\n{log[-2000:]}"
        except Exception as e: return f"讀取日誌失敗：{e}"
    return f"未知動作：{action}"


def execute_multi_monitor(action, monitor=1, window=""):
    try:
        import subprocess, win32gui, win32con
        if action == "list":
            r = subprocess.run(["powershell", "-Command", "Get-CimInstance Win32_DesktopMonitor | Select-Object Name,ScreenWidth,ScreenHeight | Format-Table -AutoSize"], capture_output=True, text=True, encoding="utf-8", errors="replace")
            return f"🖥️ 螢幕資訊：\n{r.stdout.strip()}"
        elif action in ("extend","clone"):
            mode = "/extend" if action == "extend" else "/clone"; subprocess.run(["displayswitch.exe", mode])
            return f"✅ 螢幕模式已切換為：{action}"
        elif action == "move_window":
            import ctypes; user32 = ctypes.windll.user32; sw = user32.GetSystemMetrics(0); hwnds = []
            win32gui.EnumWindows(lambda h, l: l.append(h) if win32gui.IsWindowVisible(h) and window.lower() in win32gui.GetWindowText(h).lower() else None, hwnds)
            if not hwnds: return f"⚠️ 找不到視窗：{window}"
            offset_x = sw * (int(monitor) - 1); rect = win32gui.GetWindowRect(hwnds[0])
            win32gui.MoveWindow(hwnds[0], offset_x + 100, 100, rect[2]-rect[0], rect[3]-rect[1], True)
            return f"✅ 已移動視窗 '{window}' 到螢幕 {monitor}"
    except Exception as e: return f"❌ 多螢幕管理失敗：{e}"


def execute_network_config(action, name="", ip="", dns1="", dns2="", domain="", duration=10):
    try:
        import subprocess, time
        if action == "adapter_list":
            result = subprocess.run(["powershell", "-Command",
                "Get-NetAdapter | Select-Object Name,Status,LinkSpeed | Format-Table -AutoSize"],
                capture_output=True, text=True, encoding="utf-8", errors="replace")
            return f"🌐 網路介面卡：\n{result.stdout.strip()}"
        elif action == "adapter_enable":
            subprocess.run(["powershell", "-Command", f"Enable-NetAdapter -Name '{name}' -Confirm:$false"], capture_output=True)
            return f"✅ 已啟用介面卡：{name}"
        elif action == "adapter_disable":
            subprocess.run(["powershell", "-Command", f"Disable-NetAdapter -Name '{name}' -Confirm:$false"], capture_output=True)
            return f"✅ 已停用介面卡：{name}"
        elif action == "dns_get":
            result = subprocess.run(["powershell", "-Command",
                f"Get-DnsClientServerAddress -InterfaceAlias '{name}' | Format-List"],
                capture_output=True, text=True, encoding="utf-8", errors="replace")
            return f"🌐 DNS 設定：\n{result.stdout.strip()}"
        elif action == "dns_set":
            servers = f"'{dns1}','{dns2}'" if dns2 else f"'{dns1}'"
            subprocess.run(["powershell", "-Command",
                f"Set-DnsClientServerAddress -InterfaceAlias '{name}' -ServerAddresses ({servers})"],
                capture_output=True)
            return f"✅ DNS 設定完成：{dns1}" + (f", {dns2}" if dns2 else "")
        elif action == "ip_set":
            subprocess.run(["powershell", "-Command",
                f"New-NetIPAddress -InterfaceAlias '{name}' -IPAddress '{ip}' -PrefixLength 24 -DefaultGateway '{dns1}' -Confirm:$false 2>&1"],
                capture_output=True)
            return f"✅ 靜態 IP 設定：{ip}"
        elif action == "hosts_list":
            with open(r"C:\Windows\System32\drivers\etc\hosts", "r", encoding="utf-8", errors="replace") as f:
                content = f.read()
            return f"📄 Hosts 檔案：\n{content[:1500]}"
        elif action == "hosts_add":
            entry = f"\n{ip}\t{domain}"
            with open(r"C:\Windows\System32\drivers\etc\hosts", "a", encoding="utf-8") as f:
                f.write(entry)
            return f"✅ 已新增 hosts：{ip} → {domain}"
        elif action == "hosts_remove":
            with open(r"C:\Windows\System32\drivers\etc\hosts", "r", encoding="utf-8", errors="replace") as f:
                lines = f.readlines()
            new_lines = [l for l in lines if domain not in l]
            with open(r"C:\Windows\System32\drivers\etc\hosts", "w", encoding="utf-8") as f:
                f.writelines(new_lines)
            return f"✅ 已移除 hosts：{domain}"
        elif action == "traffic":
            import psutil
            t1 = psutil.net_io_counters()
            time.sleep(min(int(duration), 10))
            t2 = psutil.net_io_counters()
            sent = (t2.bytes_sent - t1.bytes_sent) / 1024
            recv = (t2.bytes_recv - t1.bytes_recv) / 1024
            return f"📊 網路流量（{duration}s）：↑ {sent:.1f} KB  ↓ {recv:.1f} KB"
    except Exception as e:
        return f"❌ 網路設定失敗：{e}"


def execute_network_diag(action, host, ports="22,80,443,3306,3389,8080"):
    try:
        if action == "ping":
            r = subprocess.run(["ping", "-n", "4", host], capture_output=True, text=True, encoding="cp950", errors="replace", timeout=20)
            return r.stdout.strip()
        elif action == "traceroute":
            r = subprocess.run(["tracert", host], capture_output=True, text=True, encoding="cp950", errors="replace", timeout=60)
            return r.stdout[:2000]
        elif action == "portscan":
            import socket
            results = []
            for p in [int(x) for x in ports.split(",")]:
                s = socket.socket(); s.settimeout(1)
                r = s.connect_ex((host, p))
                results.append(f"Port {p}: {'開放 ✅' if r == 0 else '關閉 ❌'}")
                s.close()
            return "\n".join(results)
    except Exception as e:
        return f"❌ 網路診斷失敗：{e}"


def execute_news_monitor(action, keywords="", interval=300, duration=3600, chat_id=None, _bot_send=None):
    """全球新聞監控"""
    import feedparser, threading
    if action == "check":
        results = []; kw_list = [k.strip() for k in keywords.split(",") if k.strip()]
        for kw in kw_list[:5]:
            try:
                url = f"https://news.google.com/rss/search?q={urllib.parse.quote(kw)}&hl=zh-TW"
                feed = feedparser.parse(url)
                if feed.entries:
                    results.append(f"【{kw}】")
                    for entry in feed.entries[:3]: results.append(f"  • {entry.get('title','')}")
            except Exception: pass
        return "\n".join(results) if results else "無相關新聞"
    elif action == "top_headlines":
        try:
            cats = [("台灣", "https://news.google.com/rss?hl=zh-TW&gl=TW&ceid=TW:zh-Hant"),
                    ("科技", "https://news.google.com/rss/topics/CAAqJggKIiBDQkFTRWdvSUwyMHZNRGRqTVhZU0FtZGhLQUFQAQ?hl=zh-TW")]
            result = []
            for cat, url in cats:
                feed = feedparser.parse(url); result.append(f"【{cat}】")
                for e in feed.entries[:4]: result.append(f"  • {e.get('title','')}")
            return "\n".join(result)
        except Exception as e: return f"取得新聞失敗：{e}"
    elif action == "start_watch":
        seen = set(); kw_list = [k.strip() for k in keywords.split(",") if k.strip()]; start_t = time.time()
        def monitor():
            while time.time() - start_t < float(duration):
                for kw in kw_list:
                    try:
                        url = f"https://news.google.com/rss/search?q={urllib.parse.quote(kw)}&hl=zh-TW"; feed = feedparser.parse(url)
                        for entry in feed.entries[:3]:
                            link = entry.get("link", "")
                            if link and link not in seen:
                                seen.add(link); msg = f"🚨 新聞快報 [{kw}]\n{entry.get('title','')}"
                                if _bot_send and chat_id:
                                    import asyncio; asyncio.run_coroutine_threadsafe(_bot_send(chat_id=chat_id, text=msg), asyncio.get_event_loop())
                    except Exception: pass
                time.sleep(float(interval))
        threading.Thread(target=monitor, daemon=True).start()
        return f"✅ 開始監控：{keywords}，每 {interval}秒 檢查一次"
    return f"未知動作：{action}"


def execute_nlp(action, text):
    try:
        import anthropic as _ant
        c = _ant.Anthropic(api_key=os.getenv("ANTHROPIC_API_KEY"))
        prompt = f"請用繁體中文摘要以下文字（100字以內）：\n\n{text}" if action == "summarize" else f"分析以下文字的情緒，只回覆：正面/負面/中性 + 一句說明：\n\n{text}"
        msg = c.messages.create(model="claude-haiku-4-5-20251001", max_tokens=256, messages=[{"role":"user","content":prompt}])
        return msg.content[0].text
    except Exception as e:
        return f"❌ NLP 失敗：{e}"


def execute_notify(title, message):
    try:
        from win10toast import ToastNotifier
        ToastNotifier().show_toast(title, message, duration=5, threaded=True)
        return f"通知已送出"
    except Exception as e:
        return f"通知失敗：{e}"


def execute_ocr(image_path=""):
    import easyocr
    import numpy as np
    reader = easyocr.Reader(["ch_tra", "en"], gpu=False)
    if not image_path:
        img = pyautogui.screenshot()
        buf = io.BytesIO()
        img.save(buf, format="PNG")
        source = buf.getvalue()
    else:
        pil_img = Image.open(image_path)
        source = np.array(pil_img)
    results = reader.readtext(source)
    return "\n".join(r[1] for r in results) or "未辨識到文字"


def execute_osint_search(action, query="", target="", limit=10):
    """OSINT情報蒐集"""
    import feedparser; results = []
    if action == "web_search":
        try:
            from ddgs import DDGS
            with DDGS() as ddgs: items = list(ddgs.text(query, region="zh-tw", max_results=int(limit)))
            for r in items: results.append({"title": r.get("title",""), "url": r.get("href",""), "snippet": r.get("body","")})
        except Exception as e: return f"搜尋失敗：{e}"
        return json.dumps(results, ensure_ascii=False, indent=2) if results else "無結果"
    elif action == "news_search":
        try:
            url = f"https://news.google.com/rss/search?q={urllib.parse.quote(query)}&hl=zh-TW&gl=TW&ceid=TW:zh-Hant"
            feed = feedparser.parse(url)
            for entry in feed.entries[:int(limit)]: results.append(f"📰 {entry.get('title','')}\n   {entry.get('published','')}\n   {entry.get('link','')[:100]}")
        except Exception as e: return f"新聞搜尋失敗：{e}"
        return "\n\n".join(results) if results else "無相關新聞"
    elif action == "ip_osint":
        try:
            resp = requests.get(f"http://ip-api.com/json/{target}", timeout=5); d = resp.json()
            return f"🌐 IP情報：{target}\n國家：{d.get('country','')} ({d.get('countryCode','')})\n城市：{d.get('city','')} / {d.get('regionName','')}\nISP：{d.get('isp','')}\n組織：{d.get('org','')}\n時區：{d.get('timezone','')}\n座標：{d.get('lat','')}, {d.get('lon','')}"
        except Exception as e: return f"IP查詢失敗：{e}"
    elif action == "domain_osint":
        try:
            import socket; lines = [f"🔍 域名情報：{target}"]
            try:
                ip = socket.gethostbyname(target); lines.append(f"IP：{ip}")
                resp = requests.get(f"http://ip-api.com/json/{ip}", timeout=5); d = resp.json()
                lines.append(f"地區：{d.get('country','')} / {d.get('city','')}"); lines.append(f"ISP：{d.get('isp','')}")
            except Exception as e2: lines.append(f"DNS解析失敗：{e2}")
            return "\n".join(lines)
        except Exception as e: return f"域名查詢失敗：{e}"
    elif action == "top_news":
        try:
            feed = feedparser.parse("https://news.google.com/rss?hl=zh-TW&gl=TW&ceid=TW:zh-Hant")
            result = ["📡 Google 台灣頭條新聞："]
            for entry in feed.entries[:10]: result.append(f"• {entry.get('title','')}")
            return "\n".join(result)
        except Exception as e: return f"取得新聞失敗：{e}"
    elif action == "reddit_search":
        try:
            url = f"https://www.reddit.com/search.json?q={urllib.parse.quote(query)}&limit={limit}&sort=relevance"
            resp = requests.get(url, headers={"User-Agent": "OsintBot/1.0"}, timeout=10); data = resp.json()
            for post in data.get("data", {}).get("children", []):
                p = post["data"]; results.append(f"r/{p.get('subreddit','')} | {p.get('title','')} (👍{p.get('score',0)})\n{p.get('url','')}")
        except Exception as e: return f"Reddit搜尋失敗：{e}"
        return "\n\n".join(results) if results else "無結果"
    return f"未知動作：{action}"


def execute_password_mgr(action, site, master, username="", password=""):
    try:
        from cryptography.hazmat.primitives.kdf.pbkdf2 import PBKDF2HMAC
        from cryptography.hazmat.primitives import hashes
        from cryptography.fernet import Fernet
        import base64
        kdf = PBKDF2HMAC(algorithm=hashes.SHA256(), length=32, salt=b"pwd_manager_v1", iterations=100000)
        key = base64.urlsafe_b64encode(kdf.derive(master.encode())); f = Fernet(key)
        db = str(Path("C:/Users/blue_/claude-telegram-bot/passwords.db"))
        conn = sqlite3.connect(db)
        conn.execute("CREATE TABLE IF NOT EXISTS passwords (id INTEGER PRIMARY KEY AUTOINCREMENT, site TEXT, username TEXT, password TEXT)"); conn.commit()
        if action == "save":
            enc = f.encrypt(password.encode()).decode()
            conn.execute("INSERT INTO passwords (site,username,password) VALUES (?,?,?)", (site, username, enc)); conn.commit(); conn.close()
            return f"✅ 密碼已儲存：{site}"
        elif action == "get":
            rows = conn.execute("SELECT site,username,password FROM passwords WHERE site LIKE ?", (f"%{site}%",)).fetchall(); conn.close()
            if not rows: return f"（找不到 {site} 的密碼）"
            results = []
            for s, u, enc_p in rows:
                try:
                    decrypted = f.decrypt(enc_p.encode()).decode()
                    masked = decrypted[:2] + "*" * max(0, len(decrypted) - 4) + decrypted[-2:] if len(decrypted) > 4 else "****"
                    results.append(f"🔑 {s}\n帳號：{u}\n密碼：{masked}")
                except Exception: results.append(f"❌ {s} 解密失敗（主密碼錯誤？）")
            return "\n".join(results)
        else: conn.close(); return f"⚠️ 不支援的操作：{action}（支援 save/get）"
    except Exception as e:
        return f"❌ 密碼管理失敗：{e}"


def execute_pdf_edit(action, path="", output="", paths="", text=""):
    try:
        import fitz, json
        if action == "merge":
            pdf_list = json.loads(paths); writer = fitz.open()
            for p in pdf_list: writer.insert_pdf(fitz.open(p))
            writer.save(output)
            return f"✅ 已合併 {len(pdf_list)} 個 PDF：{output}"
        elif action == "split":
            doc = fitz.open(path); out_dir = Path(output); out_dir.mkdir(parents=True, exist_ok=True)
            for i, _ in enumerate(doc):
                out = fitz.open(); out.insert_pdf(doc, from_page=i, to_page=i)
                out.save(str(out_dir / f"page_{i+1}.pdf"))
            return f"✅ 已分割 {len(doc)} 頁到：{output}"
        elif action == "watermark":
            doc = fitz.open(path); out_path = output or path.replace(".pdf", "_wm.pdf")
            for page in doc:
                page.insert_text((page.rect.width/2-50, page.rect.height/2), text, fontsize=40, color=(0.8,0.8,0.8), rotate=45)
            doc.save(out_path)
            return f"✅ 已加浮水印：{out_path}"
    except Exception as e:
        return f"❌ PDF 編輯失敗：{e}"


def execute_pdf_to_image(path, output_dir="", dpi=150):
    try:
        import fitz
        doc = fitz.open(path)
        out = Path(output_dir) if output_dir else Path(path).parent / (Path(path).stem + "_imgs")
        out.mkdir(parents=True, exist_ok=True)
        for i, page in enumerate(doc):
            pix = page.get_pixmap(matrix=fitz.Matrix(dpi/72, dpi/72))
            pix.save(str(out / f"page_{i+1}.png"))
        return f"✅ 已轉換 {len(doc)} 頁到：{out}"
    except Exception as e:
        return f"❌ PDF 轉圖片失敗：{e}"


def execute_pentest(action, target="", port_range="1-1000", timeout=2):
    """滲透測試"""
    import socket, threading
    if action == "port_scan":
        open_ports = []
        try: target_ip = socket.gethostbyname(target)
        except Exception: target_ip = target
        if "-" in str(port_range): start, end = map(int, str(port_range).split("-")); ports = list(range(start, min(end + 1, start + 500)))
        else: ports = [int(p) for p in str(port_range).split(",")]
        def scan(port):
            try:
                s = socket.socket(); s.settimeout(float(timeout))
                if s.connect_ex((target_ip, port)) == 0:
                    try: svc = socket.getservbyport(port)
                    except Exception: svc = "unknown"
                    open_ports.append((port, svc))
                s.close()
            except Exception: pass
        threads = [threading.Thread(target=scan, args=(p,)) for p in ports]
        for t in threads: t.start()
        for t in threads: t.join(timeout=1)
        open_ports.sort()
        return f"🔍 埠掃描 {target_ip}\n開放：{len(open_ports)} 個\n" + "\n".join(f"  {p}/tcp {s}" for p, s in open_ports[:20])
    elif action == "ssl_check":
        import ssl
        try:
            ctx = ssl.create_default_context()
            with ctx.wrap_socket(socket.socket(), server_hostname=target) as s:
                s.settimeout(5); s.connect((target, 443)); cert = s.getpeercert()
            exp = dt.dt.datetime.strptime(cert['notAfter'], '%b %d %H:%M:%S %Y %Z')
            days = (exp - dt.dt.datetime.utcnow()).days
            return f"SSL憑證：{target}\n到期：{cert['notAfter']}\n剩餘：{days}天\n{'⚠️ 即將到期！' if days < 30 else '✅ 有效'}"
        except Exception as e: return f"SSL檢查失敗：{e}"
    elif action == "http_headers":
        try:
            url = f"https://{target}" if not target.startswith("http") else target
            resp = requests.get(url, timeout=5, verify=False, allow_redirects=False)
            sec_headers = ["Strict-Transport-Security", "X-Content-Type-Options", "X-Frame-Options", "Content-Security-Policy", "X-XSS-Protection", "Referrer-Policy"]
            result = [f"🔒 安全標頭：{target}"]
            for h in sec_headers: result.append(f"  {'✅' if h in resp.headers else '❌ 缺少'} {h}")
            return "\n".join(result)
        except Exception as e: return f"分析失敗：{e}"
    elif action == "vuln_scan":
        try:
            base = f"https://{target}" if not target.startswith("http") else target
            paths = [".env", ".git/config", "wp-config.php", "phpinfo.php", "admin/", "backup.sql", ".DS_Store"]
            found = []
            for path in paths:
                try:
                    r = requests.get(f"{base}/{path}", timeout=3, verify=False, allow_redirects=False)
                    if r.status_code == 200: found.append(f"  ⚠️ /{path} 可存取！")
                except Exception: pass
            if found: return f"🚨 發現 {len(found)} 個潛在漏洞：\n" + "\n".join(found)
            return f"✅ {target} 未發現常見敏感路徑暴露"
        except Exception as e: return f"掃描失敗：{e}"
    elif action == "password_audit":
        import re; results = []
        for pwd in target.split(","):
            pwd = pwd.strip(); score = 0; issues = []
            if len(pwd) >= 12: score += 2
            elif len(pwd) >= 8: score += 1
            else: issues.append("太短")
            if re.search(r'[A-Z]', pwd): score += 1
            else: issues.append("缺大寫")
            if re.search(r'[a-z]', pwd): score += 1
            else: issues.append("缺小寫")
            if re.search(r'\d', pwd): score += 1
            else: issues.append("缺數字")
            if re.search(r'[!@#$%^&*]', pwd): score += 2
            else: issues.append("缺符號")
            strength = ["極弱","弱","普通","強","極強"][min(score//2, 4)]
            results.append(f"'{pwd[:3]}***': {strength}({score}/7) {' '.join(issues)}")
        return "密碼強度審計：\n" + "\n".join(results)
    return f"未知動作：{action}"


def execute_pixel_watch(action, name="", x=0, y=0, command="", interval=1.0, tolerance=10, _bot_send=None, _chat_id=None):
    global _pixel_watchers
    try:
        import threading, time
        if action == "get":
            import pyautogui; screenshot = pyautogui.screenshot(); r, g, b = screenshot.getpixel((int(x), int(y)))[:3]
            return f"🎨 座標({x},{y}) 目前顏色：RGB({r},{g},{b}) #{r:02X}{g:02X}{b:02X}"
        elif action == "list":
            if not _pixel_watchers: return "⚠️ 無執行中的像素監控"
            return "🎨 像素監控：\n" + "\n".join(f"- {k}: ({v['x']},{v['y']})" for k,v in _pixel_watchers.items())
        elif action == "stop":
            if name in _pixel_watchers: _pixel_watchers[name]["running"] = False; del _pixel_watchers[name]; return f"✅ 已停止像素監控：{name}"
            return f"⚠️ 找不到：{name}"
        elif action == "start":
            import pyautogui; screenshot = pyautogui.screenshot(); r0, g0, b0 = screenshot.getpixel((int(x), int(y)))[:3]
            cfg = {"x": x, "y": y, "r": r0, "g": g0, "b": b0, "running": True}; _pixel_watchers[name] = cfg
            def _watch():
                import pyautogui, subprocess, time as t
                while _pixel_watchers.get(name, {}).get("running"):
                    try:
                        sc = pyautogui.screenshot(); r, g, b = sc.getpixel((int(x), int(y)))[:3]
                        cfg = _pixel_watchers.get(name, {}); diff = abs(r-cfg["r"]) + abs(g-cfg["g"]) + abs(b-cfg["b"])
                        if diff > int(tolerance) * 3:
                            msg = f"🎨 [{name}] 像素({x},{y})顏色變化！#{r:02X}{g:02X}{b:02X}"
                            if command: subprocess.Popen(command, shell=True)
                            if _bot_send and _chat_id:
                                import asyncio; asyncio.run_coroutine_threadsafe(_bot_send(_chat_id, msg), asyncio.get_event_loop())
                            cfg["r"], cfg["g"], cfg["b"] = r, g, b
                    except Exception: pass
                    t.sleep(float(interval))
            threading.Thread(target=_watch, daemon=True).start()
            return f"✅ 像素監控已啟動：{name} ({x},{y}) 容差={tolerance}"
    except Exception as e: return f"❌ 像素監控失敗：{e}"


def execute_portfolio(action: str, chat_id: int = 0, symbol: str = "",
                      shares: float = 0, cost: float = 0) -> str:
    try:
        _init_portfolio_db()
        if action == "add":
            conn = sqlite3.connect(_PORTFOLIO_DB)
            existing = conn.execute(
                "SELECT id, shares, cost FROM portfolio WHERE chat_id=? AND symbol=?",
                (chat_id, symbol.upper())
            ).fetchone()
            if existing:
                new_shares = existing[1] + shares
                new_cost = (existing[1] * existing[2] + shares * cost) / new_shares
                conn.execute("UPDATE portfolio SET shares=?, cost=? WHERE id=?",
                             (new_shares, new_cost, existing[0]))
            else:
                conn.execute(
                    "INSERT INTO portfolio (chat_id, symbol, shares, cost) VALUES (?, ?, ?, ?)",
                    (chat_id, symbol.upper(), shares, cost)
                )
            conn.commit(); conn.close()
            return f"✅ 已新增 {symbol.upper()} {shares} 股，成本 {cost:.2f}"
        elif action == "remove":
            conn = sqlite3.connect(_PORTFOLIO_DB)
            conn.execute("DELETE FROM portfolio WHERE chat_id=? AND symbol=?",
                         (chat_id, symbol.upper()))
            conn.commit(); conn.close()
            return f"✅ 已移除 {symbol.upper()}"
        elif action == "clear":
            conn = sqlite3.connect(_PORTFOLIO_DB)
            conn.execute("DELETE FROM portfolio WHERE chat_id=?", (chat_id,))
            conn.commit(); conn.close()
            return "✅ 已清空投資組合"
        elif action == "view":
            conn = sqlite3.connect(_PORTFOLIO_DB)
            rows = conn.execute(
                "SELECT symbol, shares, cost FROM portfolio WHERE chat_id=?", (chat_id,)
            ).fetchall()
            conn.close()
            if not rows:
                return "目前沒有持股紀錄，用 add 新增持股"
            import yfinance as yf
            lines = ["📁 投資組合\n"]
            total_cost = total_value = 0
            for sym, sh, ct in rows:
                try:
                    h = yf.Ticker(sym).history(period="1d")
                    cur_price = h["Close"].iloc[-1]
                    value = cur_price * sh
                    invested = ct * sh
                    pnl = value - invested
                    pnl_pct = pnl / invested * 100 if invested else 0
                    arrow = "▲" if pnl >= 0 else "▼"
                    lines.append(
                        f"📌 {sym}　{sh} 股 × {cur_price:.2f} = {value:,.0f}\n"
                        f"   成本 {ct:.2f}　損益 {arrow}{abs(pnl):,.0f} ({pnl_pct:+.1f}%)"
                    )
                    total_cost += invested; total_value += value
                except Exception:
                    lines.append(f"📌 {sym}：無法取得即時價格")
            if total_cost > 0:
                total_pnl = total_value - total_cost
                total_pct = total_pnl / total_cost * 100
                arrow = "▲" if total_pnl >= 0 else "▼"
                lines.append(
                    f"\n── 總計 ──\n"
                    f"總市值：{total_value:,.0f}\n"
                    f"總成本：{total_cost:,.0f}\n"
                    f"總損益：{arrow}{abs(total_pnl):,.0f} ({total_pct:+.1f}%)"
                )
            return "\n".join(lines)
        return "未知操作"
    except Exception as e:
        return f"投資組合操作失敗：{e}"


def execute_power(action):
    cmds = {
        "sleep": "rundll32.exe powrprof.dll,SetSuspendState 0,1,0",
        "restart": "shutdown /r /t 5",
        "shutdown": "shutdown /s /t 5",
    }
    subprocess.run(["powershell.exe", "-Command", cmds[action]])
    return f"已執行：{action}"


def execute_power_plan(action, plan="balanced"):
    try:
        import subprocess
        plan_guids = {
            "balanced": "381b4222-f694-41f0-9685-ff5bb260df2e",
            "high_performance": "8c5e7fda-e8bf-4a96-9a85-a6e23a8c635c",
            "power_saver": "a1841308-3541-4fab-bc81-f71556f20b4a",
        }
        if action == "list":
            r = subprocess.run(["powercfg", "/list"], capture_output=True, text=True, encoding="utf-8", errors="replace")
            return f"⚡ 電源計畫：\n{r.stdout.strip()}"
        elif action == "get":
            r = subprocess.run(["powercfg", "/getactivescheme"], capture_output=True, text=True, encoding="utf-8", errors="replace")
            return f"⚡ 目前電源計畫：\n{r.stdout.strip()}"
        elif action == "set":
            guid = plan_guids.get(plan)
            if not guid:
                return f"⚠️ 未知計畫：{plan}"
            subprocess.run(["powercfg", "/setactive", guid], capture_output=True)
            return f"✅ 電源計畫已切換為：{plan}"
    except Exception as e:
        return f"❌ 電源計畫失敗：{e}"


def execute_pptx(action, path, slides=""):
    try:
        from pptx import Presentation
        if action == "read":
            prs = Presentation(path); lines = []
            for i, slide in enumerate(prs.slides, 1):
                texts = [sh.text for sh in slide.shapes if sh.has_text_frame]
                lines.append(f"[投影片 {i}] " + " | ".join(t for t in texts if t.strip()))
            return "\n".join(lines) or "（簡報為空）"
        elif action == "create":
            import json; from pptx.util import Pt
            data = json.loads(slides); prs = Presentation()
            for s in data:
                sl = prs.slides.add_slide(prs.slide_layouts[1])
                if "title" in s: sl.shapes.title.text = s["title"]
                if "body" in s: sl.placeholders[1].text = s["body"]
            prs.save(path)
            return f"✅ 已建立簡報：{path}"
    except Exception as e:
        return f"❌ PPT 操作失敗：{e}"


def execute_printer(action, path="", printer_name=""):
    try:
        import subprocess, win32print
        if action == "list":
            printers = win32print.EnumPrinters(win32print.PRINTER_ENUM_LOCAL | win32print.PRINTER_ENUM_CONNECTIONS)
            default = win32print.GetDefaultPrinter()
            return f"🖨️ 印表機清單（預設：{default}）：\n" + "\n".join(f"- {p[2]}" for p in printers)
        elif action == "print":
            pname = printer_name or win32print.GetDefaultPrinter(); import win32api
            win32api.ShellExecute(0, "print", path, f'/d:"{pname}"', ".", 0)
            return f"✅ 已傳送列印：{path} → {pname}"
        elif action == "queue":
            pname = printer_name or win32print.GetDefaultPrinter(); handle = win32print.OpenPrinter(pname)
            jobs = win32print.EnumJobs(handle, 0, -1, 1); win32print.ClosePrinter(handle)
            if not jobs: return f"✅ {pname} 列印佇列為空"
            return "🖨️ 列印佇列：\n" + "\n".join(f"工作 {j['JobId']}：{j['pDocument']} ({j['Status']})" for j in jobs)
        elif action == "clear_queue":
            pname = printer_name or win32print.GetDefaultPrinter(); handle = win32print.OpenPrinter(pname)
            jobs = win32print.EnumJobs(handle, 0, -1, 1)
            for j in jobs: win32print.SetJob(handle, j["JobId"], 0, None, win32print.JOB_CONTROL_DELETE)
            win32print.ClosePrinter(handle); return f"✅ 已清除 {pname} 的列印佇列"
        elif action == "set_default": win32print.SetDefaultPrinter(printer_name); return f"✅ 預設印表機已設為：{printer_name}"
    except Exception as e: return f"❌ 印表機操作失敗：{e}"


def execute_proactive_alert(action, name="", condition="", threshold="", target="",
                             interval=60, chat_id=None, _bot_send=None):
    """主動預警系統：持續監控條件並自動通知"""
    import threading
    ALERTS_FILE = Path(__file__).parent / "proactive_alerts.json"

    def load_alerts():
        return json.loads(ALERTS_FILE.read_text("utf-8")) if ALERTS_FILE.exists() else {}

    def save_alerts(d):
        ALERTS_FILE.write_text(json.dumps(d, ensure_ascii=False, indent=2), "utf-8")

    if action == "add":
        alerts = load_alerts()
        alerts[name] = {"condition": condition, "threshold": threshold, "target": target,
                        "interval": interval, "active": True, "triggered": 0}
        save_alerts(alerts)
        return f"✅ 預警 '{name}' 建立\n條件：{condition} {threshold}\n目標：{target}\n間隔：{interval}秒"

    elif action == "list":
        alerts = load_alerts()
        if not alerts:
            return "無預警設定"
        return "🚨 預警清單：\n" + "\n".join(
            f"  {'🟢' if a.get('active') else '🔴'} [{n}] {a['condition']} {a.get('threshold','')} (觸發{a.get('triggered',0)}次)"
            for n, a in alerts.items())

    elif action == "delete":
        alerts = load_alerts()
        if name in alerts:
            del alerts[name]
            save_alerts(alerts)
            return f"✅ 預警 '{name}' 已刪除"
        return f"找不到 '{name}'"

    elif action == "toggle":
        alerts = load_alerts()
        if name in alerts:
            alerts[name]["active"] = not alerts[name].get("active", True)
            save_alerts(alerts)
            return f"✅ 預警 '{name}' {'啟用' if alerts[name]['active'] else '停用'}"
        return f"找不到 '{name}'"

    elif action == "start_all":
        alerts = load_alerts()

        def monitor_loop():
            while True:
                current = load_alerts()
                for aname, acfg in current.items():
                    if not acfg.get("active"):
                        continue
                    try:
                        msg = None
                        cond = acfg["condition"]
                        thresh = acfg.get("threshold", "")
                        tgt = acfg.get("target", "")

                        if cond == "cpu_above":
                            import psutil
                            val = psutil.cpu_percent(interval=1)
                            if val > float(thresh):
                                msg = f"🚨 CPU {val}% > {thresh}%"
                        elif cond == "memory_above":
                            import psutil
                            val = psutil.virtual_memory().percent
                            if val > float(thresh):
                                msg = f"🚨 記憶體 {val}% > {thresh}%"
                        elif cond == "price_above":
                            r = requests.get(f"https://api.binance.com/api/v3/ticker/price?symbol={tgt}", timeout=5)
                            val = float(r.json().get("price", 0))
                            if val > float(thresh):
                                msg = f"📈 {tgt} 價格 {val} > {thresh}"
                        elif cond == "price_below":
                            r = requests.get(f"https://api.binance.com/api/v3/ticker/price?symbol={tgt}", timeout=5)
                            val = float(r.json().get("price", 0))
                            if val < float(thresh):
                                msg = f"📉 {tgt} 價格 {val} < {thresh}"
                        elif cond == "news_keyword":
                            import feedparser
                            feed = feedparser.parse(f"https://news.google.com/rss/search?q={urllib.parse.quote(tgt)}&hl=zh-TW")
                            if feed.entries:
                                msg = f"📰 新聞 [{tgt}]\n{feed.entries[0].get('title','')}"

                        if msg and _bot_send and chat_id:
                            import asyncio
                            asyncio.run_coroutine_threadsafe(
                                _bot_send(chat_id=chat_id, text=f"[預警:{aname}] {msg}"),
                                asyncio.get_event_loop()
                            )
                            cur_alerts = load_alerts()
                            if aname in cur_alerts:
                                cur_alerts[aname]["triggered"] = cur_alerts[aname].get("triggered", 0) + 1
                                save_alerts(cur_alerts)
                    except Exception:
                        pass
                time.sleep(float(interval) if interval else 60)

        threading.Thread(target=monitor_loop, daemon=True).start()
        return f"✅ 已啟動 {len(alerts)} 個預警監控"

    return f"未知動作：{action}"


def execute_process_mgr(action, name="", pid=None, level="normal"):
    try:
        import psutil
        priority_map = {
            "realtime": psutil.REALTIME_PRIORITY_CLASS,
            "high": psutil.HIGH_PRIORITY_CLASS,
            "above_normal": psutil.ABOVE_NORMAL_PRIORITY_CLASS,
            "normal": psutil.NORMAL_PRIORITY_CLASS,
            "below_normal": psutil.BELOW_NORMAL_PRIORITY_CLASS,
            "idle": psutil.IDLE_PRIORITY_CLASS,
        }
        if action == "list":
            procs = sorted(psutil.process_iter(["pid","name","cpu_percent","memory_info"]), key=lambda p: p.info["cpu_percent"] or 0, reverse=True)
            lines = [f"{'PID':>6} {'CPU%':>6} {'MEM MB':>8}  名稱"]
            for p in procs[:25]:
                mem = (p.info["memory_info"].rss // 1024 // 1024) if p.info["memory_info"] else 0
                lines.append(f"{p.info['pid']:>6} {p.info['cpu_percent'] or 0:>6.1f} {mem:>8}  {p.info['name']}")
            return "\n".join(lines)
        elif action == "search":
            found = [p for p in psutil.process_iter(["pid","name","cpu_percent","memory_info"]) if name.lower() in p.info["name"].lower()]
            if not found:
                return f"⚠️ 找不到程序：{name}"
            lines = [f"PID:{p.info['pid']} CPU:{p.info['cpu_percent']}% 記憶體:{p.info['memory_info'].rss//1024//1024}MB {p.info['name']}" for p in found]
            return "\n".join(lines)
        elif action == "kill":
            targets = []
            if pid:
                targets = [psutil.Process(int(pid))]
            else:
                targets = [p for p in psutil.process_iter(["pid","name"]) if name.lower() in p.info["name"].lower()]
            if not targets:
                return f"⚠️ 找不到程序：{name}"
            for p in targets:
                p.kill()
            return f"✅ 已終止 {len(targets)} 個程序：{name or pid}"
        elif action == "priority":
            p = psutil.Process(int(pid)) if pid else next((x for x in psutil.process_iter(["pid","name"]) if name.lower() in x.info["name"].lower()), None)
            if not p:
                return f"⚠️ 找不到程序：{name}"
            p.nice(priority_map.get(level, psutil.NORMAL_PRIORITY_CLASS))
            return f"✅ 已設定 PID {p.pid} 優先權為：{level}"
    except Exception as e:
        return f"❌ 程序管理失敗：{e}"


def execute_push_notify(platform, message, webhook_or_token):
    try:
        if platform == "discord":
            resp = requests.post(webhook_or_token, json={"content": message}, timeout=10)
            return f"✅ Discord 已發送（{resp.status_code}）"
        elif platform == "line":
            resp = requests.post(
                "https://notify-api.line.me/api/notify",
                headers={"Authorization": f"Bearer {webhook_or_token}"},
                data={"message": message}, timeout=10
            )
            return f"✅ LINE 已發送（{resp.status_code}）"
    except Exception as e:
        return f"❌ 推播失敗：{e}"


def execute_qr_code(action, content="", path="", duration=30.0):
    try:
        if action == "qr_gen":
            import qrcode as _qr; img = _qr.make(content)
            save_path = path or str(Path("C:/Users/blue_/Desktop/測試檔案") / "qrcode.png"); img.save(save_path)
            return f"✅ QR Code 已生成：{save_path}"
        elif action == "qr_scan":
            from pyzbar.pyzbar import decode; from PIL import Image
            img = Image.open(path) if path else pyautogui.screenshot(); results = decode(img)
            if not results: return "❌ 未偵測到 QR Code"
            return "\n".join(f"掃描結果：{r.data.decode('utf-8')}" for r in results)
        elif action == "clipboard_watch":
            import pyperclip, time as t; last = pyperclip.paste(); changes = []; start = t.time()
            while t.time() - start < duration:
                cur = pyperclip.paste()
                if cur != last: changes.append(f"[{dt.dt.datetime.now().strftime('%H:%M:%S')}] {cur[:100]}"); last = cur
                t.sleep(0.5)
            return "\n".join(changes) if changes else f"監控 {duration} 秒內無剪貼簿變化"
        return "未知動作"
    except Exception as e: return f"操作失敗：{e}"


def execute_registry(action, key, value_name="", value=""):
    try:
        import winreg
        parts = key.split("\\", 1)
        roots = {"HKEY_LOCAL_MACHINE": winreg.HKEY_LOCAL_MACHINE,
                 "HKEY_CURRENT_USER": winreg.HKEY_CURRENT_USER,
                 "HKLM": winreg.HKEY_LOCAL_MACHINE, "HKCU": winreg.HKEY_CURRENT_USER}
        root = roots[parts[0]]
        if action == "read":
            with winreg.OpenKey(root, parts[1]) as k:
                if value_name:
                    val, _ = winreg.QueryValueEx(k, value_name)
                    return f"{value_name} = {val}"
                lines = []
                i = 0
                while True:
                    try:
                        n, v, _ = winreg.EnumValue(k, i); lines.append(f"{n} = {v}"); i += 1
                    except OSError: break
                return "\n".join(lines[:20])
        elif action == "write":
            with winreg.OpenKey(root, parts[1], 0, winreg.KEY_SET_VALUE) as k:
                winreg.SetValueEx(k, value_name, 0, winreg.REG_SZ, value)
            return f"✅ 已寫入：{value_name} = {value}"
    except Exception as e:
        return f"❌ 登錄檔操作失敗：{e}"


def execute_reminder(time_str, message):
    import threading, time as t
    def _remind():
        if time_str.isdigit(): t.sleep(int(time_str))
        else:
            import datetime as dt; now = dt.dt.datetime.now()
            target = dt.dt.datetime.strptime(time_str, "%H:%M").replace(year=now.year, month=now.month, day=now.day)
            if target < now: target = target.replace(day=now.day+1)
            t.sleep((target-now).total_seconds())
        try:
            from win10toast import ToastNotifier; ToastNotifier().show_toast("⏰ 提醒", message, duration=10)
        except Exception: pass
        try: import pyttsx3; e = pyttsx3.init(); e.say(message); e.runAndWait()
        except Exception: pass
    threading.Thread(target=_remind, daemon=True).start()
    return f"✅ 提醒已設定：{time_str} → {message}"


def execute_report(title, data_json, output=""):
    try:
        import jinja2, json
        data = json.loads(data_json)
        out_path = output or str(Path("C:/Users/blue_/Desktop/測試檔案") / f"report_{dt.dt.datetime.now().strftime('%Y%m%d_%H%M%S')}.html")
        tmpl = jinja2.Template("""<!DOCTYPE html><html><head><meta charset="utf-8">
<style>body{font-family:sans-serif;margin:40px}table{border-collapse:collapse;width:100%}
th,td{border:1px solid #ccc;padding:8px}th{background:#4472C4;color:white}
tr:nth-child(even){background:#f2f2f2}h1{color:#4472C4}</style></head>
<body><h1>{{ title }}</h1><p>生成時間：{{ time }}</p>
{% for section, rows in data.items() %}<h2>{{ section }}</h2>
{% if rows is iterable and rows is not string %}{% if rows[0] is mapping %}
<table><tr>{% for k in rows[0].keys() %}<th>{{ k }}</th>{% endfor %}</tr>
{% for row in rows %}<tr>{% for v in row.values() %}<td>{{ v }}</td>{% endfor %}</tr>{% endfor %}</table>
{% else %}<ul>{% for i in rows %}<li>{{ i }}</li>{% endfor %}</ul>{% endif %}
{% else %}<p>{{ rows }}</p>{% endif %}{% endfor %}</body></html>""")
        Path(out_path).write_text(tmpl.render(title=title, data=data, time=dt.dt.datetime.now().strftime("%Y-%m-%d %H:%M:%S")), encoding="utf-8")
        return f"✅ 報告已生成：{out_path}"
    except Exception as e:
        return f"❌ 報告生成失敗：{e}"


def execute_restore_point(action, description=""):
    try:
        if action == "create":
            ps = f"Checkpoint-Computer -Description '{description or 'Claude Auto Restore'}' -RestorePointType MODIFY_SETTINGS"
            r = subprocess.run(["powershell.exe","-Command",ps], capture_output=True, text=True, encoding="utf-8", errors="replace")
            return r.stdout or "✅ 還原點已建立"
        elif action == "list":
            r = subprocess.run(["powershell.exe","-Command","Get-ComputerRestorePoint | Select-Object SequenceNumber,Description,CreationTime | Format-Table"], capture_output=True, text=True, encoding="utf-8", errors="replace")
            return r.stdout or "（無還原點）"
    except Exception as e:
        return f"❌ 系統還原點失敗：{e}"


def execute_run_code(type_, code):
    try:
        if type_ == "python":
            import io as _io, traceback, contextlib; buf = _io.StringIO()
            with contextlib.redirect_stdout(buf):
                try: exec(compile(code, "<string>", "exec"), {})
                except Exception: buf.write(traceback.format_exc())
            return buf.getvalue() or "（執行完畢，無輸出）"
        elif type_ == "shell":
            result = subprocess.run(["powershell.exe", "-Command", code], capture_output=True, text=True, encoding="utf-8", errors="replace", timeout=30)
            return (result.stdout + result.stderr).strip() or "（執行完畢，無輸出）"
        return f"不支援的程式碼類型「{type_}」"
    except Exception as e: return f"執行失敗：{e}"


def execute_screen_record(action, duration=10.0, output=""):
    try:
        if action == "record":
            import mss, cv2, numpy as np, time as t
            out_path = output or str(Path("C:/Users/blue_/Desktop/測試檔案") / f"record_{dt.dt.datetime.now().strftime('%H%M%S')}.mp4")
            with mss.mss() as sct:
                mon = sct.monitors[1]
                w, h = mon["width"], mon["height"]
                writer = cv2.VideoWriter(out_path, cv2.VideoWriter_fourcc(*"mp4v"), 10, (w, h))
                end = t.time() + duration
                while t.time() < end:
                    frame = np.array(sct.grab(mon))
                    writer.write(cv2.cvtColor(frame, cv2.COLOR_BGRA2BGR))
                    t.sleep(0.1)
                writer.release()
            return f"✅ 錄影完成：{out_path}"
        elif action == "webcam":
            import cv2
            out_path = output or str(Path("C:/Users/blue_/Desktop/測試檔案") / f"webcam_{dt.dt.datetime.now().strftime('%H%M%S')}.jpg")
            cap = cv2.VideoCapture(0)
            ret, frame = cap.read(); cap.release()
            if ret:
                cv2.imwrite(out_path, frame)
                return f"✅ 已拍照：{out_path}"
            return "❌ 無法存取攝影機"
    except Exception as e:
        return f"❌ 失敗：{e}"


def execute_screen_vision(question: str = "請描述這個畫面上有什麼，以及目前電腦在做什麼事。") -> tuple:
    """截圖並用 Claude vision 分析，回傳 (文字分析, 截圖bytes)"""
    import base64
    img = pyautogui.screenshot()
    buf = io.BytesIO()
    img.save(buf, format="PNG")
    img_bytes = buf.getvalue()
    img_b64 = base64.standard_b64encode(img_bytes).decode("utf-8")
    response = client.messages.create(
        model="claude-sonnet-4-6",
        max_tokens=1024,
        messages=[{
            "role": "user",
            "content": [
                {"type": "image", "source": {"type": "base64", "media_type": "image/png", "data": img_b64}},
                {"type": "text", "text": question}
            ]
        }]
    )
    return response.content[0].text, img_bytes


def execute_screenshot_compare(img1_path="", img2_path="", output=""):
    try:
        import cv2, numpy as np, time as t
        img1 = cv2.cvtColor(np.array(pyautogui.screenshot()), cv2.COLOR_RGB2BGR) if not img1_path else cv2.imread(img1_path)
        if not img2_path: t.sleep(2); img2 = cv2.cvtColor(np.array(pyautogui.screenshot()), cv2.COLOR_RGB2BGR)
        else: img2 = cv2.imread(img2_path)
        h = min(img1.shape[0], img2.shape[0]); w = min(img1.shape[1], img2.shape[1])
        diff = cv2.absdiff(img1[:h,:w], img2[:h,:w])
        _, thresh = cv2.threshold(cv2.cvtColor(diff, cv2.COLOR_BGR2GRAY), 30, 255, cv2.THRESH_BINARY)
        contours, _ = cv2.findContours(thresh, cv2.RETR_EXTERNAL, cv2.CHAIN_APPROX_SIMPLE)
        result = img2[:h,:w].copy(); cv2.drawContours(result, contours, -1, (0, 0, 255), 2)
        out = output or str(Path("C:/Users/blue_/Desktop/測試檔案") / f"diff_{dt.dt.datetime.now().strftime('%H%M%S')}.png")
        cv2.imwrite(out, result)
        pct = cv2.countNonZero(thresh) / (h * w) * 100
        return f"差異：{pct:.2f}%，標記圖：{out}"
    except Exception as e:
        return f"❌ 截圖比對失敗：{e}"


def execute_self_benchmark(action):
    """自我評估：測試各功能健康狀態並產生報告"""
    if action == "run":
        tests = {
            "網路連線": lambda: requests.get("https://google.com", timeout=3).status_code == 200,
            "Claude API Key": lambda: bool(os.environ.get("ANTHROPIC_API_KEY", "")),
            "Telegram Token": lambda: bool(os.environ.get("BOT_TOKEN", "")),
            "資料庫": lambda: sqlite3.connect(Path(__file__).parent / "memory.db").execute("SELECT 1").fetchone() is not None,
            "截圖功能": lambda: bool(__import__("pyautogui").screenshot()),
            "HuggingFace Key": lambda: bool(os.environ.get("HF_TOKEN", "")),
            "知識庫": lambda: (Path(__file__).parent / "knowledge_base.db").exists(),
            "目標管理": lambda: (Path(__file__).parent / "goals.db").exists(),
            "Binance Key": lambda: bool(os.environ.get("BINANCE_KEY", "")),
            "HA Token": lambda: bool(os.environ.get("HA_TOKEN", "")),
            "VirusTotal Key": lambda: bool(os.environ.get("VIRUSTOTAL_KEY", "")),
            "預警設定": lambda: (Path(__file__).parent / "proactive_alerts.json").exists(),
        }
        passed = 0
        lines = ["🤖 自我評估報告\n"]
        for name, test in tests.items():
            try:
                ok = test()
                lines.append(f"  {'✅' if ok else '⚠️'} {name}")
                if ok:
                    passed += 1
            except Exception as e:
                lines.append(f"  ❌ {name}: {str(e)[:30]}")
        score = int(passed / len(tests) * 100)
        lines.insert(1, f"健康度：{score}% ({passed}/{len(tests)})")
        return "\n".join(lines)

    elif action == "skill_count":
        import re
        content = Path(__file__).read_text(encoding="utf-8")
        skills = sorted(set(re.findall(r'def execute_(\w+)\(', content)))
        return f"技能函數總數：{len(skills)}\n\n" + "\n".join(f"• {s}" for s in skills)

    elif action == "memory_stats":
        db = Path(__file__).parent / "memory.db"
        if not db.exists():
            return "資料庫不存在"
        conn = sqlite3.connect(db)
        hist = conn.execute("SELECT COUNT(*) FROM chat_history").fetchone()[0]
        memo = conn.execute("SELECT COUNT(*) FROM long_term_memory").fetchone()[0]
        conn.close()
        kb_db = Path(__file__).parent / "knowledge_base.db"
        kb_count = 0
        if kb_db.exists():
            c2 = sqlite3.connect(kb_db)
            kb_count = c2.execute("SELECT COUNT(*) FROM knowledge").fetchone()[0]
            c2.close()
        return f"記憶統計：\n對話歷史：{hist} 條\n長期記憶：{memo} 條\n知識庫：{kb_count} 條"

    return f"未知動作：{action}"


def execute_send_email(to, subject, body):
    import smtplib
    from email.mime.text import MIMEText
    from email.mime.multipart import MIMEMultipart
    smtp_user = os.getenv("SMTP_USER", "")
    smtp_pass = os.getenv("SMTP_PASS", "")
    if not smtp_user: return "未設定 SMTP_USER / SMTP_PASS，請加入 .env"
    msg = MIMEMultipart()
    msg["From"] = smtp_user; msg["To"] = to; msg["Subject"] = subject
    msg.attach(MIMEText(body, "plain", "utf-8"))
    with smtplib.SMTP(os.getenv("SMTP_HOST","smtp.gmail.com"), int(os.getenv("SMTP_PORT","587"))) as s:
        s.starttls(); s.login(smtp_user, smtp_pass); s.send_message(msg)
    return f"Email 已寄出到 {to}"


def execute_serial_port(action, port="", baudrate=9600, data="", timeout=2):
    try:
        import serial, serial.tools.list_ports
        if action == "list":
            ports = serial.tools.list_ports.comports()
            if not ports: return "❌ 未找到 COM port"
            return "\n".join(f"{p.device} - {p.description}" for p in ports)
        if not port: return "❌ 需指定 port"
        if action == "send":
            with serial.Serial(port, baudrate, timeout=timeout) as s: s.write(data.encode())
            return f"✅ 已發送：{data}"
        elif action == "read":
            with serial.Serial(port, baudrate, timeout=timeout) as s: resp = s.read(1024).decode("utf-8", errors="ignore")
            return f"收到：{resp}" if resp else "❌ 無資料回應"
        elif action == "send_read":
            with serial.Serial(port, baudrate, timeout=timeout) as s:
                s.write(data.encode()); import time; time.sleep(0.1); resp = s.read(1024).decode("utf-8", errors="ignore")
            return f"發送：{data}\n收到：{resp}"
        return "未知動作"
    except ImportError: return "❌ 請先安裝：pip install pyserial"
    except Exception as e: return f"❌ serial_port 失敗：{e}"


def execute_smart_home(action, device="", value="", host="", token=""):
    """智慧家居控制：Home Assistant API整合"""
    ha_host = host or os.environ.get("HA_HOST", "http://homeassistant.local:8123")
    ha_token = token or os.environ.get("HA_TOKEN", "")
    headers = {"Authorization": f"Bearer {ha_token}", "Content-Type": "application/json"}

    if action == "list_devices":
        try:
            resp = requests.get(f"{ha_host}/api/states", headers=headers, timeout=10)
            if resp.status_code == 200:
                states = resp.json()
                result = [f"🏠 智慧家居設備（{len(states)} 個）："]
                for s in states[:25]:
                    result.append(f"  {s['entity_id']}: {s['state']}")
                return "\n".join(result)
            return f"連線失敗 {resp.status_code} — 請設定 HA_HOST 和 HA_TOKEN"
        except Exception as e:
            return f"連線失敗：{e}\n請設定環境變數 HA_HOST 和 HA_TOKEN"

    elif action in ("turn_on", "turn_off"):
        try:
            domain = device.split(".")[0] if "." in device else "homeassistant"
            svc = "turn_on" if action == "turn_on" else "turn_off"
            resp = requests.post(f"{ha_host}/api/services/{domain}/{svc}",
                                 headers=headers, json={"entity_id": device}, timeout=10)
            return f"✅ {device} {'開啟' if action=='turn_on' else '關閉'}" if resp.status_code == 200 else f"失敗：{resp.status_code}"
        except Exception as e:
            return f"控制失敗：{e}"

    elif action == "get_state":
        try:
            resp = requests.get(f"{ha_host}/api/states/{device}", headers=headers, timeout=10)
            if resp.status_code == 200:
                s = resp.json()
                return f"{device}\n狀態：{s['state']}\n屬性：{json.dumps(s.get('attributes',{}), ensure_ascii=False)[:300]}"
            return f"查詢失敗：{resp.status_code}"
        except Exception as e:
            return f"查詢失敗：{e}"

    elif action == "set_value":
        try:
            domain = device.split(".")[0] if "." in device else "input_number"
            resp = requests.post(f"{ha_host}/api/services/{domain}/set_value",
                                 headers=headers, json={"entity_id": device, "value": value}, timeout=10)
            return f"✅ {device} 設定為 {value}" if resp.status_code == 200 else f"失敗：{resp.status_code}"
        except Exception as e:
            return f"設定失敗：{e}"

    elif action == "run_scene":
        try:
            resp = requests.post(f"{ha_host}/api/services/scene/turn_on",
                                 headers=headers, json={"entity_id": f"scene.{device}"}, timeout=10)
            return f"✅ 場景 {device} 已啟動" if resp.status_code == 200 else f"失敗：{resp.status_code}"
        except Exception as e:
            return f"執行失敗：{e}"

    return f"未知動作：{action}"


def execute_software(action, name="", keyword=""):
    try:
        import subprocess
        if action == "list":
            q = f"| Where-Object {{$_.DisplayName -like '*{keyword}*'}}" if keyword else ""
            result = subprocess.run(["powershell", "-Command",
                f"Get-ItemProperty HKLM:\\Software\\Wow6432Node\\Microsoft\\Windows\\CurrentVersion\\Uninstall\\* {q} | Select-Object DisplayName,DisplayVersion | Sort-Object DisplayName | Format-Table -AutoSize"],
                capture_output=True, text=True, encoding="utf-8", errors="replace")
            return f"📦 已安裝軟體：\n{result.stdout.strip()[:2000]}"
        elif action == "install":
            result = subprocess.run(["winget", "install", "--id", name, "-e", "--silent"],
                capture_output=True, text=True)
            if result.returncode == 0:
                return f"✅ 安裝成功：{name}"
            return f"⚠️ 安裝結果：{result.stdout[-500:]}"
        elif action == "uninstall":
            result = subprocess.run(["winget", "uninstall", "--name", name, "--silent"],
                capture_output=True, text=True)
            if result.returncode == 0:
                return f"✅ 已卸載：{name}"
            return f"⚠️ 卸載結果：{result.stdout[-500:]}"
    except Exception as e:
        return f"❌ 軟體管理失敗：{e}"


def execute_sound_detect(action, threshold=20, duration=5, output=""):
    try:
        import numpy as _np
        try: import sounddevice as _sd
        except ImportError: return "❌ 請先安裝：pip install sounddevice"
        RATE = 16000
        if action == "volume_level":
            data = _sd.rec(int(RATE * 1), samplerate=RATE, channels=1, dtype="float32"); _sd.wait()
            vol = int(_np.abs(data).mean() * 1000); return f"🔊 當前音量：{vol}/100"
        elif action == "detect_silence":
            import time as _t; silent_start = None; start = _t.time()
            while _t.time() - start < duration:
                data = _sd.rec(int(RATE * 0.5), samplerate=RATE, channels=1, dtype="float32"); _sd.wait()
                vol = int(_np.abs(data).mean() * 1000)
                if vol < threshold:
                    if silent_start is None: silent_start = _t.time()
                    elif _t.time() - silent_start > 1.5: return f"🔇 偵測到靜音（音量 {vol}）"
                else: silent_start = None
            return f"監控 {duration} 秒內無靜音段"
        elif action == "detect_speech":
            import time as _t; start = _t.time()
            while _t.time() - start < duration:
                data = _sd.rec(int(RATE * 0.5), samplerate=RATE, channels=1, dtype="float32"); _sd.wait()
                vol = int(_np.abs(data).mean() * 1000)
                if vol > threshold: return f"🗣 偵測到說話（音量 {vol}）"
            return f"監控 {duration} 秒內未偵測到說話"
        elif action == "record_until_silence":
            import time as _t, wave; frames = []
            out = output or str(__import__("pathlib").Path("C:/Users/blue_/Desktop/測試檔案") / "recording.wav")
            CHUNK = int(RATE * 0.5); silent_count = 0
            while True:
                data = _sd.rec(CHUNK, samplerate=RATE, channels=1, dtype="int16"); _sd.wait()
                frames.append(data.tobytes()); vol = int(_np.abs(data.astype(float)).mean() / 32768 * 1000)
                if vol < threshold: silent_count += 1
                else: silent_count = 0
                if silent_count >= 4 or len(frames) > 200: break
            with wave.open(out, "wb") as wf: wf.setnchannels(1); wf.setsampwidth(2); wf.setframerate(RATE); wf.writeframes(b"".join(frames))
            return f"✅ 錄音完成：{out}（{len(frames)*0.5:.1f}秒）"
        return "未知動作"
    except Exception as e: return f"❌ sound_detect 失敗：{e}"


def execute_speedtest():
    try:
        import speedtest as st
        s = st.Speedtest(); s.get_best_server()
        dl = s.download()/1_000_000; ul = s.upload()/1_000_000
        return f"📶 下載：{dl:.1f} Mbps | 上傳：{ul:.1f} Mbps | Ping：{s.results.ping:.0f} ms | 伺服器：{s.results.server.get('name','')}"
    except Exception as e:
        return f"❌ 速度測試失敗：{e}"


def execute_ssh_sftp(action, host, user, password, command="", local="", remote="", port=22):
    try:
        import paramiko
        if action == "ssh_run":
            c = paramiko.SSHClient(); c.set_missing_host_key_policy(paramiko.AutoAddPolicy())
            c.connect(host, port=port, username=user, password=password, timeout=15)
            _, stdout, stderr = c.exec_command(command)
            out = stdout.read().decode(errors="replace") + stderr.read().decode(errors="replace"); c.close()
            return out.strip() or "（執行完畢，無輸出）"
        else:
            t = paramiko.Transport((host, port)); t.connect(username=user, password=password)
            sftp = paramiko.SFTPClient.from_transport(t)
            if action == "sftp_upload": sftp.put(local, remote); result = f"✅ 已上傳：{local} → {remote}"
            else: sftp.get(remote, local); result = f"✅ 已下載：{remote} → {local}"
            sftp.close(); t.close(); return result
    except Exception as e: return f"❌ SSH/SFTP 失敗：{e}"


def execute_startup(action, name="", command=""):
    try:
        import subprocess, winreg
        key_path = r"Software\Microsoft\Windows\CurrentVersion\Run"
        if action == "list":
            result = subprocess.run(["powershell", "-Command",
                "Get-CimInstance Win32_StartupCommand | Select-Object Name,Command,Location | Format-Table -AutoSize"],
                capture_output=True, text=True, encoding="utf-8", errors="replace")
            return f"🚀 開機自啟：\n{result.stdout.strip()}"
        elif action == "add":
            with winreg.OpenKey(winreg.HKEY_CURRENT_USER, key_path, 0, winreg.KEY_SET_VALUE) as k:
                winreg.SetValueEx(k, name, 0, winreg.REG_SZ, command)
            return f"✅ 已新增開機自啟：{name}"
        elif action == "remove":
            try:
                with winreg.OpenKey(winreg.HKEY_CURRENT_USER, key_path, 0, winreg.KEY_SET_VALUE) as k:
                    winreg.DeleteValue(k, name)
                return f"✅ 已移除開機自啟：{name}"
            except FileNotFoundError:
                return f"⚠️ 找不到項目：{name}"
    except Exception as e:
        return f"❌ 開機自啟管理失敗：{e}"


def execute_stt(duration=5):
    import sounddevice as sd
    import soundfile as sf
    import speech_recognition as sr
    import tempfile
    sample_rate = 16000
    recording = sd.rec(int(duration * sample_rate), samplerate=sample_rate, channels=1, dtype="int16")
    sd.wait()
    with tempfile.NamedTemporaryFile(suffix=".wav", delete=False) as f:
        tmp_path = f.name
    sf.write(tmp_path, recording, sample_rate)
    r = sr.Recognizer()
    with sr.AudioFile(tmp_path) as source:
        audio = r.record(source)
    Path(tmp_path).unlink(missing_ok=True)
    try:
        return r.recognize_google(audio, language="zh-TW")
    except Exception as e:
        return f"語音辨識失敗：{e}"


def execute_system_monitor(action, target=""):
    try:
        if action == "info":
            cpu = psutil.cpu_percent(interval=1)
            mem = psutil.virtual_memory()
            disk = psutil.disk_usage("C:/")
            return (f"CPU：{cpu}%\n"
                    f"記憶體：{mem.percent}%（{mem.used//1024//1024}MB / {mem.total//1024//1024}MB）\n"
                    f"磁碟 C：{disk.percent}%（{disk.used//1024//1024//1024}GB / {disk.total//1024//1024//1024}GB）")
        elif action == "process_list":
            procs = sorted(psutil.process_iter(["pid","name","memory_info"]),
                           key=lambda p: p.info["memory_info"].rss if p.info["memory_info"] else 0, reverse=True)
            lines = [f"PID:{p.info['pid']} {p.info['name']} ({p.info['memory_info'].rss//1024//1024}MB)"
                     for p in procs[:20] if p.info["memory_info"]]
            return "\n".join(lines)
        elif action == "kill":
            try:
                psutil.Process(int(target)).kill()
                return f"已結束 PID {target}"
            except ValueError:
                killed = sum(1 for p in psutil.process_iter(["name"]) if target.lower() in p.info["name"].lower() and not p.kill())
                return f"已結束 {killed} 個「{target}」"
    except Exception as e:
        return f"執行失敗：{e}"


def execute_system_tools(action, **kwargs):
    try:
        if action == "event_log":
            log_name = kwargs.get("log_name","System"); level = kwargs.get("level","Error"); count = kwargs.get("count",10)
            r = subprocess.run(["powershell.exe","-Command",f"Get-WinEvent -LogName '{log_name}' -MaxEvents {count} | Where-Object {{$_.LevelDisplayName -eq '{level}'}} | Select-Object TimeCreated,Message | Format-List"], capture_output=True, text=True, encoding="utf-8", errors="replace", timeout=30)
            return r.stdout[:2000] or f"（無 {level} 事件）"
        elif action == "usb_list":
            r = subprocess.run(["powershell.exe","-Command","Get-PnpDevice | Where-Object {$_.Class -eq 'USB' -and $_.Status -eq 'OK'} | Select-Object FriendlyName | Format-Table"], capture_output=True, text=True, encoding="utf-8", errors="replace")
            return r.stdout[:2000] or "（無 USB 裝置）"
        elif action == "firewall_list":
            r = subprocess.run(["powershell.exe","-Command","Get-NetFirewallRule | Where-Object {$_.Enabled -eq 'True'} | Select-Object DisplayName,Direction,Action | Format-Table -AutoSize"], capture_output=True, text=True, encoding="utf-8", errors="replace")
            return r.stdout[:2000]
        elif action == "firewall_add":
            name=kwargs.get("name",""); direction=kwargs.get("direction","in"); port=kwargs.get("port",80)
            r = subprocess.run(["powershell.exe","-Command",f"New-NetFirewallRule -DisplayName '{name}' -Direction {direction} -Protocol TCP -LocalPort {port} -Action Allow"], capture_output=True, text=True, encoding="utf-8", errors="replace")
            return r.stdout or f"✅ 防火牆規則已新增：{name}"
        elif action == "firewall_remove":
            r = subprocess.run(["powershell.exe","-Command",f"Remove-NetFirewallRule -DisplayName '{kwargs.get('name','')}'"], capture_output=True, text=True, encoding="utf-8", errors="replace")
            return r.stdout or f"✅ 防火牆規則已刪除"
        elif action == "printer_list":
            r = subprocess.run(["powershell.exe","-Command","Get-Printer | Select-Object Name,PrinterStatus | Format-Table"], capture_output=True, text=True, encoding="utf-8", errors="replace")
            return r.stdout or "（無印表機）"
        elif action == "font_list":
            kw = kwargs.get("keyword","")
            r = subprocess.run(["powershell.exe","-Command","[System.Reflection.Assembly]::LoadWithPartialName('System.Drawing') | Out-Null; [System.Drawing.FontFamily]::Families | Select-Object -ExpandProperty Name"], capture_output=True, text=True, encoding="utf-8", errors="replace")
            fonts = [f for f in r.stdout.strip().splitlines() if kw.lower() in f.lower()] if kw else r.stdout.strip().splitlines()
            return "\n".join(fonts[:50]) + (f"\n...共 {len(fonts)} 個" if len(fonts)>50 else "")
        elif action == "rdp": subprocess.Popen(["mstsc", f"/v:{kwargs.get('host','')}"]); return f"✅ 正在開啟 RDP：{kwargs.get('host','')}"
    except Exception as e: return f"❌ 失敗：{e}"


# ============================================================
# GPU OCR 腳本互斥鎖（防止同時跑兩個 GPU 腳本導致 VRAM 爆炸）
# ============================================================
_gpu_script_lock = __import__("threading").Lock()
_gpu_script_active = None


def _check_gpu_available(script_name):
    """檢查 GPU 是否被其他腳本佔用"""
    global _gpu_script_active, _tg_auto_reply_proc
    with _gpu_script_lock:
        if _gpu_script_active == "tg_auto_reply" and _tg_auto_reply_proc is not None:
            try:
                if _tg_auto_reply_proc.poll() is None:
                    return False, f"GPU 被 Telegram 自動回覆佔用中，請先停止後再執行 {script_name}"
            except Exception:
                pass
            _gpu_script_active = None
        return True, ""


def _set_gpu_active(script_name):
    global _gpu_script_active
    with _gpu_script_lock:
        _gpu_script_active = script_name


def _release_gpu(script_name):
    global _gpu_script_active
    with _gpu_script_lock:
        if _gpu_script_active == script_name:
            _gpu_script_active = None


_tg_auto_reply_proc = None
_TG_AUTO_STOP_FILE = "C:/Users/blue_/Desktop/測試檔案/.stop_auto_reply"
_TG_AUTO_SCRIPT = "C:/Users/blue_/claude-telegram-bot/scripts/tg_auto_chat.py"


def execute_tg_auto_reply(action: str = "start", duration_minutes: float = 30, stop_time: str = "", contact_name: str = "") -> str:
    """開啟/停止 Telegram 自動回覆監控（subprocess 呼叫 tg_auto_chat.py）"""
    global _tg_auto_reply_proc

    if action == "stop":
        try:
            with open(_TG_AUTO_STOP_FILE, "w") as f:
                f.write("stop")
        except Exception:
            pass
        _tg_auto_reply_proc = None
        _release_gpu("tg_auto_reply")
        return "自動回覆已停止"

    if _tg_auto_reply_proc is not None:
        try:
            if _tg_auto_reply_proc.poll() is None:
                return "自動回覆已在運行中"
        except Exception:
            pass
        _tg_auto_reply_proc = None

    if stop_time:
        end_str = stop_time
    else:
        end_dt = dt.dt.datetime.now() + dt.timedelta(minutes=duration_minutes)
        end_str = end_dt.strftime("%H:%M")

    if not contact_name:
        return "請提供要自動回覆的好友名稱（contact_name 參數）"

    # GPU 互斥檢查
    available, reason = _check_gpu_available("tg_auto_reply")
    if not available:
        return reason

    try:
        _tg_auto_reply_proc = subprocess.Popen(
            [sys.executable, _TG_AUTO_SCRIPT, contact_name, end_str],
            cwd="C:/Users/blue_/claude-telegram-bot",
        )
        _set_gpu_active("tg_auto_reply")
        return f"自動回覆已開啟：對象 {contact_name}，監控到 {end_str}"
    except Exception as e:
        return f"自動回覆啟動失敗：{e}"


def execute_think_as(person: str, question: str, list_available: bool = False) -> str:
    """載入蒸餾好的人物思維框架，用該人物角度分析問題"""
    from pathlib import Path
    skills_dir = Path(__file__).parent / "skills"

    if list_available:
        available = [f.stem for f in skills_dir.glob("*.md") if f.stem != "colleague-niuma"]
        return "🧠 可用的思維框架：\n" + "\n".join(f"- {name}" for name in available) if available else "目前沒有可用的思維框架"

    # 找對應的 skill 檔案
    slug = person.lower().strip().replace(" ", "-").replace("_", "-")
    # 別名對照
    aliases = {
        "馬斯克": "elon-musk", "musk": "elon-musk", "elon": "elon-musk",
        "巴菲特": "warren-buffett", "buffett": "warren-buffett", "warren": "warren-buffett",
        "黃仁勳": "jensen-huang", "jensen": "jensen-huang", "huang": "jensen-huang",
        "張忠謀": "morris-chang", "morris": "morris-chang", "chang": "morris-chang",
    }
    slug = aliases.get(slug, slug)
    skill_path = skills_dir / f"{slug}.md"

    if not skill_path.exists():
        available = [f.stem for f in skills_dir.glob("*.md") if f.stem != "colleague-niuma"]
        return f"找不到「{person}」的思維框架。可用的有：{', '.join(available)}"

    # 讀取 skill 內容
    skill_content = skill_path.read_text(encoding="utf-8")

    # 用 Claude 載入框架分析問題
    from anthropic import Anthropic
    client = Anthropic()
    resp = client.messages.create(
        model="claude-sonnet-4-6",
        max_tokens=1500,
        system=f"你現在要完全用以下人物的思維框架來分析問題。不是模仿說話方式，是用他的心智模型和決策規則來推導。\n\n{skill_content}",
        messages=[{"role": "user", "content": question}]
    )
    return f"🧠 {person} 的視角：\n\n{resp.content[0].text}"


def execute_threat_intel(action, target="", api_key=""):
    """威脅情報查詢"""
    if action in ("check_url", "check_hash", "check_ip"):
        vt_key = api_key or os.environ.get("VIRUSTOTAL_KEY", "")
        if not vt_key: return "需要設定環境變數 VIRUSTOTAL_KEY"
        headers = {"x-apikey": vt_key}
        try:
            if action == "check_url":
                import base64; url_id = base64.urlsafe_b64encode(target.encode()).decode().strip("=")
                resp = requests.get(f"https://www.virustotal.com/api/v3/urls/{url_id}", headers=headers, timeout=10)
            elif action == "check_ip": resp = requests.get(f"https://www.virustotal.com/api/v3/ip_addresses/{target}", headers=headers, timeout=10)
            else: resp = requests.get(f"https://www.virustotal.com/api/v3/files/{target}", headers=headers, timeout=10)
            if resp.status_code == 200:
                stats = resp.json().get("data", {}).get("attributes", {}).get("last_analysis_stats", {})
                malicious = stats.get("malicious", 0); suspicious = stats.get("suspicious", 0); total = sum(stats.values()) if stats else 0
                level = "🔴 高危" if malicious > 3 else ("🟡 可疑" if malicious > 0 else "🟢 安全")
                return f"威脅分析：{target}\n狀態：{level}\n惡意：{malicious}/{total}\n可疑：{suspicious}/{total}"
            return f"查詢失敗：HTTP {resp.status_code}"
        except Exception as e: return f"查詢失敗：{e}"
    elif action == "scan_connections":
        try:
            result = subprocess.run(["netstat", "-ano"], capture_output=True, text=True, timeout=10)
            lines = [l for l in result.stdout.split("\n") if "ESTABLISHED" in l]
            external = [l for l in lines if not any(ip in l for ip in ["127.0.0.1", "0.0.0.0", "::1"])]
            return f"活躍外部連線（{len(external)} 條）：\n" + "\n".join(external[:15])
        except Exception as e: return f"掃描失敗：{e}"
    elif action == "check_abuse_ip":
        key = api_key or os.environ.get("ABUSEIPDB_KEY", "")
        if not key: return "需要設定環境變數 ABUSEIPDB_KEY"
        try:
            resp = requests.get("https://api.abuseipdb.com/api/v2/check", params={"ipAddress": target, "maxAgeInDays": 90},
                headers={"Key": key, "Accept": "application/json"}, timeout=10)
            if resp.status_code == 200:
                d = resp.json().get("data", {})
                return f"IP威脅：{target}\n濫用信心度：{d.get('abuseConfidenceScore',0)}%\n舉報次數：{d.get('totalReports',0)}\n國家：{d.get('countryCode','')}\nISP：{d.get('isp','')}"
            return f"查詢失敗：{resp.status_code}"
        except Exception as e: return f"查詢失敗：{e}"
    return f"未知動作：{action}"


def execute_todo(action, task="", todo_id=0):
    try:
        db = str(Path("C:/Users/blue_/claude-telegram-bot/todo.db"))
        conn = sqlite3.connect(db)
        conn.execute("CREATE TABLE IF NOT EXISTS todos (id INTEGER PRIMARY KEY AUTOINCREMENT, task TEXT, done INTEGER DEFAULT 0, created TEXT)")
        conn.commit()
        if action == "add": conn.execute("INSERT INTO todos (task,created) VALUES (?,?)", (task, dt.dt.datetime.now().strftime("%Y-%m-%d %H:%M"))); conn.commit(); conn.close(); return f"✅ 已新增：{task}"
        elif action == "list":
            rows = conn.execute("SELECT id,task,done,created FROM todos ORDER BY done,id").fetchall(); conn.close()
            return "\n".join(f"{'✅' if r[2] else '⬜'} [{r[0]}] {r[1]}" for r in rows) or "（清單為空）"
        elif action == "done": conn.execute("UPDATE todos SET done=1 WHERE id=?", (todo_id,)); conn.commit(); conn.close(); return f"✅ 任務 #{todo_id} 已完成"
        elif action == "delete": conn.execute("DELETE FROM todos WHERE id=?", (todo_id,)); conn.commit(); conn.close(); return f"✅ 任務 #{todo_id} 已刪除"
        elif action == "clear": conn.execute("DELETE FROM todos WHERE done=1"); conn.commit(); conn.close(); return "✅ 已清除所有已完成任務"
        else: conn.close(); return f"⚠️ 不支援的操作：{action}"
    except Exception as e:
        return f"❌ 任務清單失敗：{e}"


def execute_translate(text, target="zh-TW", source="auto"):
    try:
        from deep_translator import GoogleTranslator; return GoogleTranslator(source=source, target=target).translate(text)
    except Exception as e: return f"❌ 翻譯失敗：{e}"


def execute_tts(text):
    import pyttsx3
    engine = pyttsx3.init()
    engine.setProperty("rate", 180)
    engine.say(clean_for_tts(text))
    engine.runAndWait()
    return f"已朗讀完畢"


def execute_tts_advanced(action, text="", voice="zh-CN-YunxiNeural"):
    try:
        import edge_tts, asyncio
        if action == "speak":
            out = str(Path("C:/Users/blue_/Desktop/測試檔案") / f"tts_{dt.dt.datetime.now().strftime('%H%M%S')}.mp3")
            _clean = clean_for_tts(text)
            async def _gen():
                comm = edge_tts.Communicate(_clean, voice, rate="-5%", pitch="-5Hz")
                await comm.save(out)
            asyncio.run(_gen())
            subprocess.Popen(["powershell.exe","-Command",f"Start-Process '{out}'"])
            return f"✅ Edge TTS 語音已播放：{voice}"
        elif action == "list_voices":
            async def _list(): return await edge_tts.list_voices()
            voices = asyncio.run(_list())
            return "\n".join(f"{v['ShortName']}  {v['FriendlyName']}" for v in voices if v["Locale"].startswith("zh"))
    except Exception as e:
        return f"❌ Edge TTS 失敗：{e}"


def execute_ui_auto(action, window="", control="", text=""):
    try:
        from pywinauto import Desktop, Application; desktop = Desktop(backend="uia")
        if action == "get_windows":
            wins = [str(w.window_text()) for w in desktop.windows() if w.window_text()]
            return "🪟 所有視窗：\n" + "\n".join(f"- {w}" for w in wins[:30])
        win = None
        for w in desktop.windows():
            if window.lower() in w.window_text().lower(): win = w; break
        if not win: return f"⚠️ 找不到視窗：{window}"
        if action == "find":
            ctrls = win.descendants(); info = [f"[{c.control_type()}] {c.window_text()}" for c in ctrls if c.window_text()][:30]
            return f"🔍 視窗 '{window}' 的控制項：\n" + "\n".join(info)
        elif action == "read":
            texts = [c.window_text() for c in win.descendants() if c.window_text()]
            return f"📖 '{window}' 內容：\n" + "\n".join(texts[:50])
        elif action == "click":
            for c in win.descendants():
                if control.lower() in c.window_text().lower(): c.click_input(); return f"✅ 已點擊：{c.window_text()}"
            return f"⚠️ 找不到控制項：{control}"
        elif action == "type":
            for c in win.descendants():
                if control.lower() in c.window_text().lower() or c.control_type() in ("Edit","Document"):
                    c.type_keys(text); return f"✅ 已輸入文字到：{c.window_text() or c.control_type()}"
            return f"⚠️ 找不到輸入框：{control}"
    except Exception as e: return f"❌ UI 自動化失敗：{e}"


def execute_user_account(action, username="", password=""):
    try:
        import subprocess
        if action == "list":
            result = subprocess.run(["powershell", "-Command",
                "Get-LocalUser | Select-Object Name,Enabled,LastLogon | Format-Table -AutoSize"],
                capture_output=True, text=True, encoding="utf-8", errors="replace")
            return f"👤 使用者帳戶：\n{result.stdout.strip()}"
        elif action == "create":
            ps = f"$pw = ConvertTo-SecureString '{password}' -AsPlainText -Force; New-LocalUser '{username}' -Password $pw -FullName '{username}'"
            result = subprocess.run(["powershell", "-Command", ps], capture_output=True, text=True)
            if result.returncode == 0:
                return f"✅ 已建立帳戶：{username}"
            return f"❌ 建立失敗：{result.stderr.strip()}"
        elif action == "delete":
            result = subprocess.run(["powershell", "-Command", f"Remove-LocalUser '{username}'"],
                capture_output=True, text=True)
            if result.returncode == 0:
                return f"✅ 已刪除帳戶：{username}"
            return f"❌ 刪除失敗：{result.stderr.strip()}"
    except Exception as e:
        return f"❌ 使用者帳戶管理失敗：{e}"


def execute_vdesktop(action):
    acts = {"left": ("ctrl","win","left"), "right": ("ctrl","win","right"), "new": ("ctrl","win","d")}
    if action not in acts:
        return f"不支援的操作「{action}」，支援：{', '.join(acts.keys())}"
    pyautogui.hotkey(*acts[action])
    return f"虛擬桌面：{action}"


def execute_video_gen(mode: str = "slideshow", output: str = "", **kwargs) -> str:
    """
    生成影片
    mode: slideshow / text_video / tts_video / screen_record
    """
    import re as _re
    import numpy as np
    import subprocess
    import tempfile
    import traceback
    from pathlib import Path
    from PIL import Image, ImageDraw, ImageFont
    import imageio_ffmpeg

    ffmpeg_exe = imageio_ffmpeg.get_ffmpeg_exe()
    w, h = kwargs.get("size", (1280, 720))
    fps  = kwargs.get("fps", 24)
    out  = output or str(Path("C:/Users/blue_/Desktop/測試檔案") / f"video_{dt.dt.datetime.now().strftime('%H%M%S')}.mp4")

    def _write_frames(frames_iter, out_path, vid_fps, width, height):
        proc = subprocess.Popen([
            ffmpeg_exe, "-y",
            "-f", "rawvideo", "-vcodec", "rawvideo",
            "-s", f"{width}x{height}", "-pix_fmt", "rgb24",
            "-r", str(vid_fps), "-i", "pipe:0",
            "-vcodec", "libx264", "-pix_fmt", "yuv420p",
            "-preset", "fast", "-crf", "23",
            out_path
        ], stdin=subprocess.PIPE, stderr=subprocess.PIPE)
        for frame in frames_iter:
            proc.stdin.write(np.asarray(frame, dtype=np.uint8).tobytes())
        proc.stdin.close()
        rc  = proc.wait()
        err = proc.stderr.read().decode(errors="replace")
        if rc != 0 or not Path(out_path).exists():
            raise RuntimeError(f"ffmpeg 失敗 rc={rc}: {err[-300:]}")

    def _get_font(pt=60):
        for fp in ["C:/Windows/Fonts/msjh.ttc", "C:/Windows/Fonts/msyh.ttc",
                   "C:/Windows/Fonts/simhei.ttf", "C:/Windows/Fonts/arial.ttf"]:
            if Path(fp).exists():
                try:
                    return ImageFont.truetype(fp, pt)
                except Exception:
                    pass
        return ImageFont.load_default()

    def _tts_sync(text: str, voice: str, out_path: str):
        import asyncio
        import edge_tts
        loop = asyncio.new_event_loop()
        asyncio.set_event_loop(loop)
        try:
            comm = edge_tts.Communicate(text, voice, rate="-5%", pitch="-5Hz")
            loop.run_until_complete(comm.save(out_path))
        finally:
            loop.close()
        if not Path(out_path).exists():
            raise RuntimeError("TTS 語音檔案未生成")

    def _get_audio_dur(audio_path: str) -> float:
        r = subprocess.run([ffmpeg_exe, "-i", audio_path],
                           capture_output=True, text=True, errors="replace")
        m = _re.search(r"Duration:\s*(\d+):(\d+):(\d+\.?\d*)", r.stderr)
        if m:
            return int(m.group(1)) * 3600 + int(m.group(2)) * 60 + float(m.group(3))
        return 10.0

    try:
        if mode == "slideshow":
            images   = kwargs.get("images", [])
            dur      = kwargs.get("duration", 3)
            sl_fps   = kwargs.get("fps", 12)
            trans    = kwargs.get("transition", 0.5)
            if not images:
                return "❌ 請提供 images 參數（圖片路徑列表）"
            sf = int(sl_fps * dur)
            tf = int(sl_fps * trans)
            def _frames():
                loaded = []
                for p in images:
                    try:
                        loaded.append(np.array(Image.open(p).convert("RGB").resize((w, h), Image.LANCZOS)))
                    except Exception:
                        loaded.append(np.zeros((h, w, 3), dtype=np.uint8))
                for arr in loaded:
                    for i in range(tf):
                        yield (arr * (i / tf)).astype(np.uint8)
                    for _ in range(max(1, sf - tf * 2)):
                        yield arr
                    for i in range(tf):
                        yield (arr * (1.0 - i / tf)).astype(np.uint8)
            _write_frames(_frames(), out, sl_fps, w, h)
            return f"✅ 投影片影片已生成：{out}"

        elif mode == "text_video":
            text     = kwargs.get("text", "Hello")
            dur      = kwargs.get("duration", 5)
            bg_col   = tuple(kwargs.get("bg_color",   [30, 30, 40]))
            fg_col   = tuple(kwargs.get("font_color", [255, 255, 255]))
            fsize    = kwargs.get("font_size", 60)
            font     = _get_font(fsize)
            total    = max(1, int(fps * dur))
            def _frames():
                for i in range(total):
                    p = i / total
                    a = p / 0.2 if p < 0.2 else (1.0 - p) / 0.2 if p > 0.8 else 1.0
                    img  = Image.new("RGB", (w, h), bg_col)
                    draw = ImageDraw.Draw(img)
                    c    = tuple(int(x * a) for x in fg_col)
                    lines = text.split("\n")
                    lh    = fsize + 10
                    y0    = (h - len(lines) * lh) // 2
                    for j, ln in enumerate(lines):
                        bb = draw.textbbox((0, 0), ln, font=font)
                        draw.text(((w - (bb[2] - bb[0])) // 2, y0 + j * lh), ln, font=font, fill=c)
                    yield np.array(img)
            _write_frames(_frames(), out, fps, w, h)
            return f"✅ 文字動畫影片已生成：{out}"

        elif mode == "tts_video":
            text      = kwargs.get("text", "")
            img_path  = kwargs.get("image", "")
            voice     = kwargs.get("voice", "zh-CN-YunxiNeural")
            subtitle  = kwargs.get("subtitle", True)
            if not text:
                return "❌ 請提供 text 參數"
            tmp     = Path(tempfile.mkdtemp())
            audio   = str(tmp / "tts.mp3")
            vidtmp  = str(tmp / "silent.mp4")
            _tts_sync(clean_for_tts(text), voice, audio)
            dur   = _get_audio_dur(audio)
            total = max(1, int(fps * dur))
            if img_path and Path(img_path).exists():
                bg = np.array(Image.open(img_path).convert("RGB").resize((w, h), Image.LANCZOS))
            else:
                bg = np.zeros((h, w, 3), dtype=np.uint8)
                for row in range(h):
                    v = int(20 + 30 * row / h)
                    bg[row, :] = [v, v, v + 20]
            font  = _get_font(40)
            cpl   = 20
            subs  = [text[i:i+cpl] for i in range(0, len(text), cpl)]
            def _frames():
                for i in range(total):
                    frame = bg.copy()
                    if subtitle and subs:
                        ln  = subs[min(int(i / total * len(subs)), len(subs) - 1)]
                        img = Image.fromarray(frame)
                        d   = ImageDraw.Draw(img)
                        bb  = d.textbbox((0, 0), ln, font=font)
                        tw  = bb[2] - bb[0]
                        tx  = (w - tw) // 2
                        ty  = int(h * 0.82)
                        d.rectangle([tx - 10, ty - 5, tx + tw + 10, ty + 45], fill=(0, 0, 0))
                        d.text((tx, ty), ln, font=font, fill=(255, 255, 255))
                        frame = np.array(img)
                    yield frame
            _write_frames(_frames(), vidtmp, fps, w, h)
            r = subprocess.run([
                ffmpeg_exe, "-y", "-i", vidtmp, "-i", audio,
                "-c:v", "copy", "-c:a", "aac", "-shortest", out
            ], capture_output=True)
            if not Path(out).exists():
                raise RuntimeError(f"音訊合併失敗：{r.stderr.decode(errors='replace')[-200:]}")
            return f"✅ TTS 語音影片已生成：{out}"

        elif mode == "screen_record":
            import mss, time as _t
            dur      = kwargs.get("duration", 10)
            rec_fps  = kwargs.get("fps", 10)
            interval = 1.0 / rec_fps
            total    = max(1, int(rec_fps * dur))
            with mss.mss() as sct:
                mon    = sct.monitors[1]
                sw, sh = mon["width"], mon["height"]
                def _frames():
                    for _ in range(total):
                        t0  = _t.time()
                        arr = np.array(sct.grab(mon))[:, :, :3][:, :, ::-1]
                        yield arr
                        elapsed = _t.time() - t0
                        if elapsed < interval:
                            _t.sleep(interval - elapsed)
                _write_frames(_frames(), out, rec_fps, sw, sh)
            return f"✅ 螢幕錄影完成：{out}（{dur} 秒）"
        else:
            return f"❌ 未知 mode：{mode}"
    except Exception as e:
        return f"❌ 影片生成失敗：{e}\n{traceback.format_exc()}"


def execute_video_gif(path, start=0, duration=5.0, output="", fps=10):
    try:
        import imageio
        out = output or path.replace(".mp4", ".gif")
        reader = imageio.get_reader(path)
        video_fps = reader.get_meta_data().get("fps", 30)
        frames = [f for i, f in enumerate(reader) if int(start*video_fps) <= i < int((start+duration)*video_fps)]
        imageio.mimsave(out, frames, fps=fps)
        return f"✅ GIF 已生成：{out}（{len(frames)} 幀）"
    except Exception as e:
        return f"❌ 影片轉 GIF 失敗：{e}"


def execute_video_process(action, path, second=0, start=0, end=0, output=""):
    try:
        import cv2
        if action == "screenshot":
            cap = cv2.VideoCapture(path)
            cap.set(cv2.CAP_PROP_POS_MSEC, second * 1000)
            ret, frame = cap.read()
            cap.release()
            out = output or path.replace(".mp4", f"_frame{int(second)}s.jpg")
            if ret:
                cv2.imwrite(out, frame)
                return f"✅ 已擷取畫面：{out}"
            return "❌ 無法讀取影片"
        elif action == "trim":
            out = output or path.replace(".mp4", "_trim.mp4")
            subprocess.run(["ffmpeg", "-y", "-i", path, "-ss", str(start), "-to", str(end), "-c", "copy", out], capture_output=True)
            return f"✅ 已剪輯：{out}"
    except Exception as e:
        return f"❌ 影片處理失敗：{e}"


def execute_vision_loop(goal, max_steps=20, interval=3.0, timeout=120.0):
    try:
        import pyautogui, anthropic, base64, io, time
        steps = 0; start = time.time(); log = []
        _client = anthropic.Anthropic(api_key=os.getenv("ANTHROPIC_API_KEY"))
        while steps < int(max_steps) and (time.time() - start) < float(timeout):
            screenshot = pyautogui.screenshot(); buf = io.BytesIO()
            screenshot.save(buf, format="PNG")
            img_b64 = base64.standard_b64encode(buf.getvalue()).decode()
            resp = _client.messages.create(model="claude-sonnet-4-6", max_tokens=512,
                messages=[{"role": "user", "content": [
                    {"type": "image", "source": {"type": "base64", "media_type": "image/png", "data": img_b64}},
                    {"type": "text", "text": f"目標：{goal}\n已執行步驟：{log}\n\n請分析畫面，回答：\n1. 目標是否已達成？（是/否）\n2. 如果否，下一步應該怎麼做？請用 JSON 格式回答：{{\"done\": true/false, \"action\": \"動作說明\", \"type\": \"click/type/key/wait\", \"x\": 0, \"y\": 0, \"text\": \"\"}}"}
                ]}])
            import json, re
            text = resp.content[0].text; m = re.search(r'\{.*\}', text, re.DOTALL)
            if not m: log.append(f"步驟{steps+1}: AI回應無法解析"); steps += 1; time.sleep(float(interval)); continue
            action = json.loads(m.group())
            if action.get("done"): return f"✅ 目標達成！共執行 {steps} 步\n" + "\n".join(log)
            act_type = action.get("type",""); act_desc = action.get("action","")
            if act_type == "click" and action.get("x"): pyautogui.click(action["x"], action["y"])
            elif act_type == "type" and action.get("text"): pyautogui.typewrite(action["text"], interval=0.05)
            elif act_type == "key" and action.get("text"): pyautogui.press(action["text"])
            elif act_type == "wait": time.sleep(2)
            log.append(f"步驟{steps+1}: {act_desc}"); steps += 1; time.sleep(float(interval))
        return f"⏳ 達到上限（{steps} 步 / {int(time.time()-start)}s）\n執行記錄：\n" + "\n".join(log)
    except Exception as e:
        return f"❌ 視覺自動化循環失敗：{e}"


def execute_voice_cmd(action, duration=300.0, language="zh-TW", _bot_send=None, _chat_id=None):
    global _voice_cmd_running
    try:
        import threading
        if action == "stop": _voice_cmd_running = False; return "✅ 語音命令模式已停止"
        elif action == "start":
            if _voice_cmd_running: return "⚠️ 語音命令模式已在執行中"
            _voice_cmd_running = True
            def _listen_loop():
                global _voice_cmd_running
                import sounddevice as sd, soundfile as sf, speech_recognition as sr, tempfile, time, subprocess
                recognizer = sr.Recognizer(); sample_rate = 16000; end_time = time.time() + float(duration)
                while _voice_cmd_running and time.time() < end_time:
                    try:
                        recording = sd.rec(int(4 * sample_rate), samplerate=sample_rate, channels=1, dtype="int16"); sd.wait()
                        with tempfile.NamedTemporaryFile(suffix=".wav", delete=False) as f: tmp = f.name
                        sf.write(tmp, recording, sample_rate)
                        with sr.AudioFile(tmp) as source: audio = recognizer.record(source)
                        Path(tmp).unlink(missing_ok=True)
                        text = recognizer.recognize_google(audio, language=language)
                        if not text: continue
                        if "停止" in text or "stop" in text.lower(): _voice_cmd_running = False; break
                        subprocess.Popen(text, shell=True)
                    except Exception: pass
                _voice_cmd_running = False
            threading.Thread(target=_listen_loop, daemon=True).start()
            return f"✅ 語音命令模式已啟動（{duration}s），說「停止」結束"
    except Exception as e:
        return f"❌ 語音命令模式失敗：{e}"


def execute_voice_id(action, name="", audio_path="", duration=5):
    """聲紋辨識"""
    VOICE_DIR = Path(__file__).parent / "voice_profiles"; VOICE_DIR.mkdir(exist_ok=True)
    META_FILE = Path(__file__).parent / "voice_profiles.json"
    def load_meta(): return json.loads(META_FILE.read_text("utf-8")) if META_FILE.exists() else {}
    def save_meta(d): META_FILE.write_text(json.dumps(d, ensure_ascii=False, indent=2), "utf-8")
    if action == "enroll":
        try:
            import sounddevice as sd, scipy.io.wavfile as wf, numpy as np
            sr = 16000; rec = sd.rec(int(float(duration) * sr), samplerate=sr, channels=1); sd.wait()
            p = VOICE_DIR / f"{name}.wav"; wf.write(str(p), sr, rec)
            meta = load_meta(); meta[name] = str(p); save_meta(meta)
            return f"✅ {name} 聲紋已登記"
        except ImportError: return "需要安裝：pip install sounddevice scipy"
        except Exception as e: return f"登記失敗：{e}"
    elif action == "identify":
        try:
            import librosa, numpy as np; meta = load_meta()
            if not meta: return "尚未登記任何聲紋"
            if not audio_path: return "請提供音訊檔案路徑"
            q_audio, _ = librosa.load(audio_path, sr=16000); q_mfcc = librosa.feature.mfcc(y=q_audio, sr=16000, n_mfcc=13).mean(axis=1)
            best, best_score = None, float('inf')
            for person, path in meta.items():
                p_audio, _ = librosa.load(path, sr=16000); p_mfcc = librosa.feature.mfcc(y=p_audio, sr=16000, n_mfcc=13).mean(axis=1)
                score = np.linalg.norm(q_mfcc - p_mfcc)
                if score < best_score: best_score, best = score, person
            conf = max(0, 100 - int(best_score * 10))
            return f"聲紋辨識：{best}（信心度 {conf}%）"
        except ImportError: return "需要安裝：pip install librosa"
        except Exception as e: return f"辨識失敗：{e}"
    elif action == "list":
        meta = load_meta()
        return "已登記聲紋：\n" + "\n".join(f"• {n}" for n in meta) if meta else "尚未登記任何聲紋"
    elif action == "delete":
        meta = load_meta()
        if name in meta: Path(meta[name]).unlink(missing_ok=True); del meta[name]; save_meta(meta); return f"✅ {name} 聲紋已刪除"
        return f"找不到 {name}"
    return f"未知動作：{action}"


def execute_volume(action, level=None):
    try:
        from ctypes import cast, POINTER
        from comtypes import CLSCTX_ALL
        from pycaw.pycaw import AudioUtilities, IAudioEndpointVolume
        devices = AudioUtilities.GetSpeakers()
        interface = devices.Activate(IAudioEndpointVolume._iid_, CLSCTX_ALL, None)
        volume = cast(interface, POINTER(IAudioEndpointVolume))
        if action == "get":
            vol = int(volume.GetMasterVolumeLevelScalar() * 100)
            muted = bool(volume.GetMute())
            return f"🔊 音量：{vol}%{'（靜音中）' if muted else ''}"
        elif action == "set":
            volume.SetMasterVolumeLevelScalar(max(0, min(100, int(level))) / 100.0, None)
            return f"✅ 音量設定為 {level}%"
        elif action == "mute":
            volume.SetMute(1, None)
            return "✅ 已靜音"
        elif action == "unmute":
            volume.SetMute(0, None)
            return "✅ 已取消靜音"
    except Exception as e:
        return f"❌ 音量控制失敗：{e}"


def execute_vpn(action, name="", user="", password=""):
    try:
        if action == "list":
            r = subprocess.run(["powershell.exe","-Command","Get-VpnConnection | Select-Object Name,ConnectionStatus | Format-Table"], capture_output=True, text=True, encoding="utf-8", errors="replace")
            return r.stdout or "（未設定 VPN）"
        elif action == "connect":
            r = subprocess.run(["rasdial", name, user, password], capture_output=True, text=True, encoding="cp950", errors="replace")
            return r.stdout.strip()
        elif action == "disconnect":
            r = subprocess.run(["rasdial", name, "/disconnect"], capture_output=True, text=True, encoding="cp950", errors="replace")
            return r.stdout.strip()
    except Exception as e:
        return f"❌ VPN 失敗：{e}"


def execute_wait_for_text(text, timeout=60.0, interval=2.0, region=""):
    try:
        import pyautogui, easyocr, time, tempfile
        reader = easyocr.Reader(["ch_tra","en"], gpu=False); start = time.time()
        reg = None
        if region:
            parts = [int(v) for v in region.split(",")]
            if len(parts) == 4: reg = tuple(parts)
        while time.time() - start < float(timeout):
            screenshot = pyautogui.screenshot(region=reg)
            tmp = tempfile.mktemp(suffix=".png"); screenshot.save(tmp)
            results = reader.readtext(tmp, detail=0); Path(tmp).unlink(missing_ok=True)
            full = " ".join(results)
            if text.lower() in full.lower():
                elapsed = time.time() - start
                return f"✅ 偵測到文字「{text}」（等待 {elapsed:.1f}s）"
            time.sleep(float(interval))
        return f"⏳ 超時（{timeout}s），未偵測到文字「{text}」"
    except Exception as e:
        return f"❌ 等待文字失敗：{e}"


def execute_wait_seconds(seconds):
    try:
        import time
        s = float(seconds)
        time.sleep(min(s, 60))
        return f"✅ 已等待 {s} 秒"
    except Exception as e:
        return f"❌ 等待失敗：{e}"


def execute_wake_listen(keyword="小牛馬", duration=5):
    try:
        import sounddevice as sd, soundfile as sf, speech_recognition as sr, tempfile, time
        sample_rate = 16000; deadline = time.time() + 60
        while time.time() < deadline:
            recording = sd.rec(int(int(duration) * sample_rate), samplerate=sample_rate, channels=1, dtype="int16"); sd.wait()
            with tempfile.NamedTemporaryFile(suffix=".wav", delete=False) as f: tmp_path = f.name
            sf.write(tmp_path, recording, sample_rate); recognizer = sr.Recognizer()
            try:
                with sr.AudioFile(tmp_path) as source: audio = recognizer.record(source)
                text = recognizer.recognize_google(audio, language="zh-TW")
                Path(tmp_path).unlink(missing_ok=True)
                if keyword in text: return f"✅ 偵測到喚醒詞！辨識到：{text}"
            except Exception: Path(tmp_path).unlink(missing_ok=True)
        return f"⏳ 60 秒內未偵測到喚醒詞「{keyword}」"
    except Exception as e:
        return f"❌ 語音監聽失敗：{e}"


def execute_wake_word(action, keyword="", duration=5, language="zh-TW"):
    try:
        import speech_recognition as _sr
        r = _sr.Recognizer(); mic = _sr.Microphone()
        if action == "listen_once":
            with mic as src: r.adjust_for_ambient_noise(src, duration=0.5); audio = r.listen(src, timeout=duration, phrase_time_limit=duration)
            try: text = r.recognize_google(audio, language=language); return f"🎤 聽到：{text}"
            except _sr.UnknownValueError: return "❌ 無法辨識語音"
        elif action == "transcribe_stream":
            import time as _t; results = []; end = _t.time() + duration
            while _t.time() < end:
                with mic as src:
                    r.adjust_for_ambient_noise(src, duration=0.3)
                    try: audio = r.listen(src, timeout=2, phrase_time_limit=5); text = r.recognize_google(audio, language=language); results.append(text)
                    except (_sr.WaitTimeoutError, _sr.UnknownValueError): pass
            return "\n".join(results) if results else "未偵測到語音"
        elif action == "detect_keyword":
            import time as _t; end = _t.time() + duration
            while _t.time() < end:
                with mic as src:
                    r.adjust_for_ambient_noise(src, duration=0.3)
                    try:
                        audio = r.listen(src, timeout=2, phrase_time_limit=5); text = r.recognize_google(audio, language=language)
                        if keyword.lower() in text.lower(): return f"✅ 偵測到關鍵字「{keyword}」，完整語音：{text}"
                    except (_sr.WaitTimeoutError, _sr.UnknownValueError): pass
            return f"❌ {duration} 秒內未偵測到關鍵字「{keyword}」"
        return "未知動作"
    except ImportError: return "❌ 請先安裝：pip install SpeechRecognition"
    except Exception as e: return f"❌ wake_word 失敗：{e}"


def execute_watchdog(process, script, duration=60.0):
    import psutil, time as t
    restarts = 0
    end = t.time() + duration
    while t.time() < end:
        running = any(p.name().lower() == process.lower() for p in psutil.process_iter())
        if not running:
            subprocess.Popen(["pythonw", script] if script.endswith(".py") else [script])
            restarts += 1
        t.sleep(5)
    return f"守護結束，共重啟 {restarts} 次"


def execute_web_monitor(action, url, selector="body", interval=60, duration=300, keyword=""):
    try:
        import requests as _req, time as _t, hashlib, json
        from pathlib import Path as _P
        try: from bs4 import BeautifulSoup as _BS
        except ImportError: _BS = None
        CACHE_FILE = str(_P.home() / ".web_monitor_cache.json")
        def _fetch(u, sel):
            r = _req.get(u, timeout=15, headers={"User-Agent": "Mozilla/5.0"}); r.raise_for_status()
            if _BS:
                soup = _BS(r.text, "html.parser"); el = soup.select_one(sel)
                return el.get_text(strip=True)[:2000] if el else soup.get_text(strip=True)[:2000]
            return r.text[:2000]
        if action == "check_once": return f"📄 {url}\n{_fetch(url, selector)}"
        elif action == "get_price":
            text = _fetch(url, selector); import re
            prices = re.findall(r"[\$NT\$￥]?\s*[\d,]+\.?\d*", text)
            return f"💰 找到的價格：{', '.join(prices[:10])}" if prices else "❌ 未找到價格"
        elif action == "diff":
            try:
                with open(CACHE_FILE, encoding="utf-8") as f: cache = json.load(f)
            except Exception: cache = {}
            current = _fetch(url, selector); prev = cache.get(url, "")
            cache[url] = current
            with open(CACHE_FILE, "w", encoding="utf-8") as f: json.dump(cache, f, ensure_ascii=False)
            if not prev: return f"✅ 已記錄初始狀態（{len(current)} 字）"
            if prev == current: return "✅ 與上次相同，無變化"
            import difflib
            diff = list(difflib.unified_diff(prev.splitlines(), current.splitlines(), lineterm="", n=2))
            return "⚠️ 內容有變化：\n" + "\n".join(diff[:30])
        elif action == "watch":
            last = _fetch(url, selector); changes = []; start = _t.time()
            while _t.time() - start < duration:
                _t.sleep(interval)
                try:
                    current = _fetch(url, selector)
                    if current != last:
                        msg = f"⚠️ [{_t.strftime('%H:%M:%S')}] 網頁內容變化"
                        if keyword:
                            if keyword in current and keyword not in last: msg += f"（出現關鍵字「{keyword}」）"
                            elif keyword not in current and keyword in last: msg += f"（關鍵字「{keyword}」消失）"
                        changes.append(msg); last = current
                except Exception as ex: changes.append(f"❌ 抓取失敗：{ex}")
            return "\n".join(changes) if changes else f"監控 {duration} 秒內無變化"
        return "未知動作"
    except Exception as e:
        return f"❌ web_monitor 失敗：{e}"


def execute_web_scrape(action, url="", selector="body", interval=2.0, region="full"):
    try:
        if action == "scrape":
            from bs4 import BeautifulSoup
            res = requests.get(url, timeout=15, headers={"User-Agent": "Mozilla/5.0"})
            soup = BeautifulSoup(res.text, "html.parser")
            static_text = soup.get_text(strip=True)
            use_playwright = len(static_text) < 500 or len(soup.find_all("script")) > 10
            if use_playwright:
                try:
                    import subprocess, sys, json as _json, textwrap
                    _script = textwrap.dedent(f"""
import sys
from playwright.sync_api import sync_playwright
with sync_playwright() as p:
    b = p.chromium.launch(headless=True)
    page = b.new_page(user_agent="Mozilla/5.0 (Windows NT 10.0; Win64; x64) AppleWebKit/537.36")
    page.goto({_json.dumps(url)}, wait_until="networkidle", timeout=20000)
    page.wait_for_timeout(2000)
    els = page.query_selector_all({_json.dumps(selector)})
    texts = [e.inner_text() for e in els[:10] if e.inner_text().strip()]
    b.close()
    print("\\n".join(texts))
""")
                    proc = subprocess.run([sys.executable, "-c", _script], capture_output=True, encoding="utf-8", timeout=30,
                        env={**__import__("os").environ, "PYTHONIOENCODING": "utf-8"})
                    if proc.returncode == 0 and proc.stdout.strip(): return proc.stdout.strip()
                except Exception: pass
            elements = soup.select(selector)
            static_result = "\n".join(e.get_text(strip=True) for e in elements[:10])
            return static_result if static_result else f"（網頁內容太少，可能為 JS 動態網站：{url}）"
        elif action == "screen_diff":
            import time as t; img1 = pyautogui.screenshot(); t.sleep(interval); img2 = pyautogui.screenshot()
            import numpy as np; a1, a2 = np.array(img1), np.array(img2)
            diff = np.abs(a1.astype(int) - a2.astype(int)).mean()
            if diff > 2.0: return f"⚠️ 螢幕有變化（差異度：{diff:.1f}）"
            return f"✅ 螢幕無明顯變化（差異度：{diff:.1f}）"
        return "未知動作"
    except Exception as e:
        return f"操作失敗：{e}"


def execute_webcam(action, duration=5.0, output="", device=0):
    try:
        import cv2, tempfile
        if action == "list":
            found = []
            for i in range(5):
                cap = cv2.VideoCapture(i)
                if cap.isOpened():
                    found.append(f"裝置 {i}")
                    cap.release()
            return f"📷 可用攝影機：\n" + ("\n".join(found) if found else "無")
        elif action == "photo":
            cap = cv2.VideoCapture(int(device))
            if not cap.isOpened(): return f"❌ 無法開啟攝影機 {device}"
            ret, frame = cap.read(); cap.release()
            if not ret: return "❌ 無法拍攝"
            out = output or str(Path("C:/Users/blue_/Desktop/測試檔案") / f"webcam_{dt.dt.datetime.now().strftime('%Y%m%d_%H%M%S')}.jpg")
            cv2.imwrite(out, frame)
            return f"✅ 已拍照：{out}"
        elif action == "video":
            cap = cv2.VideoCapture(int(device))
            if not cap.isOpened(): return f"❌ 無法開啟攝影機 {device}"
            out = output or str(Path("C:/Users/blue_/Desktop/測試檔案") / f"webcam_{dt.dt.datetime.now().strftime('%Y%m%d_%H%M%S')}.avi")
            w = int(cap.get(cv2.CAP_PROP_FRAME_WIDTH))
            h = int(cap.get(cv2.CAP_PROP_FRAME_HEIGHT))
            writer = cv2.VideoWriter(out, cv2.VideoWriter_fourcc(*"XVID"), 20, (w,h))
            import time
            end = time.time() + float(duration)
            while time.time() < end:
                ret, frame = cap.read()
                if ret: writer.write(frame)
            cap.release(); writer.release()
            return f"✅ 已錄影 {duration}s：{out}"
    except Exception as e:
        return f"❌ 攝影機操作失敗：{e}"


def execute_webhook_server(action, port=8765, secret=""):
    global _webhook_server_proc
    try:
        import socket, subprocess as _sp, sys, os as _os
        if action == "start":
            if _webhook_server_proc and _webhook_server_proc.poll() is None:
                return f"✅ Webhook 伺服器已在運行（port {port}）"
            script = f"""
import http.server, json, threading, time
log = []
class H(http.server.BaseHTTPRequestHandler):
    def do_POST(self):
        length = int(self.headers.get('Content-Length', 0))
        body = self.rfile.read(length).decode('utf-8', errors='ignore')
        secret = "{secret}"
        if secret and self.headers.get('X-Secret','') != secret:
            self.send_response(403); self.end_headers(); return
        log.append(f"[{{time.strftime('%H:%M:%S')}}] {{self.path}}: {{body[:200]}}")
        open(r'C:/Users/blue_/claude-telegram-bot/webhook_log.txt','a',encoding='utf-8').write(log[-1]+'\\n')
        self.send_response(200); self.send_header('Content-Type','text/plain'); self.end_headers()
        self.wfile.write(b'ok')
    def log_message(self, *a): pass
http.server.HTTPServer(('0.0.0.0', {port}), H).serve_forever()
"""
            _webhook_server_proc = _sp.Popen([sys.executable, "-c", script], creationflags=0x00000008)
            return f"✅ Webhook 伺服器已啟動 port {port}\n本機 IP：{socket.gethostbyname(socket.gethostname())}"
        elif action == "stop":
            if _webhook_server_proc: _webhook_server_proc.terminate(); _webhook_server_proc = None
            return "✅ Webhook 伺服器已停止"
        elif action == "status":
            running = _webhook_server_proc and _webhook_server_proc.poll() is None
            log_path = r"C:/Users/blue_/claude-telegram-bot/webhook_log.txt"
            log = ""
            if _os.path.exists(log_path):
                with open(log_path, encoding="utf-8") as f: log = "".join(f.readlines()[-5:])
            return f"狀態：{'運行中' if running else '已停止'}\n最近事件：\n{log}"
        elif action == "get_url":
            import socket
            ip = socket.gethostbyname(socket.gethostname())
            return f"Webhook URL：http://{ip}:{port}/\n（需同一區域網路）"
        return "未知動作"
    except Exception as e:
        return f"❌ webhook_server 失敗：{e}"


def execute_webpage_shot(action, url, selector="body", interval=60.0, duration=3600.0):
    try:
        if action == "screenshot":
            from playwright.sync_api import sync_playwright
            out = str(Path("C:/Users/blue_/Desktop/測試檔案") / f"webpage_{dt.dt.datetime.now().strftime('%H%M%S')}.png")
            with sync_playwright() as p:
                browser = p.chromium.launch(); page = browser.new_page(viewport={"width": 1280, "height": 800})
                page.goto(url, timeout=30000); page.wait_for_load_state("networkidle")
                page.screenshot(path=out, full_page=True); browser.close()
            return f"✅ 網頁截圖已存：{out}"
        elif action == "monitor":
            import hashlib, time as t
            from bs4 import BeautifulSoup
            def _fetch():
                r = requests.get(url, timeout=15, headers={"User-Agent": "Mozilla/5.0"})
                soup = BeautifulSoup(r.text, "html.parser")
                text = "\n".join(e.get_text(strip=True) for e in soup.select(selector))
                return hashlib.md5(text.encode()).hexdigest(), text[:150]
            last_hash, _ = _fetch(); end = t.time() + duration; changes = []
            while t.time() < end:
                t.sleep(interval)
                new_hash, snippet = _fetch()
                if new_hash != last_hash:
                    changes.append(f"[{dt.dt.datetime.now().strftime('%H:%M:%S')}] {snippet}")
                    last_hash = new_hash
            return f"監控結束，共 {len(changes)} 次變化\n" + "\n".join(changes[:5])
    except Exception as e:
        return f"❌ 失敗：{e}"


def execute_win_notify_relay(action, duration=3600.0, filter_app="", _bot_send=None, _chat_id=None):
    try:
        import threading, time, subprocess
        if action == "status": return "ℹ️ Windows 通知攔截透過輪詢事件記錄實現"
        elif action == "stop": return "✅ 通知攔截已標記停止（重啟生效）"
        elif action == "start":
            def _relay():
                import win32evtlog, win32evtlogutil, time as t
                seen = set(); end = t.time() + float(duration); hand = win32evtlog.OpenEventLog(None, "Application")
                while t.time() < end:
                    try:
                        flags = win32evtlog.EVENTLOG_BACKWARDS_READ | win32evtlog.EVENTLOG_SEQUENTIAL_READ
                        batch = win32evtlog.ReadEventLog(hand, flags, 0)
                        for e in (batch or []):
                            eid = (e.RecordNumber, e.TimeGenerated.Format())
                            if eid in seen: continue
                            seen.add(eid); src = e.SourceName
                            if filter_app and filter_app.lower() not in src.lower(): continue
                            try: msg = win32evtlogutil.SafeFormatMessage(e, "Application")[:200]
                            except Exception: msg = "(無法讀取)"
                            if _bot_send and _chat_id and e.EventType in (1, 2):
                                import asyncio; level = "❌" if e.EventType == 1 else "⚠️"
                                asyncio.run_coroutine_threadsafe(_bot_send(_chat_id, f"🔔 Windows 通知\n{level} {src}\n{msg}"), asyncio.get_event_loop())
                    except Exception: pass
                    t.sleep(10)
                win32evtlog.CloseEventLog(hand)
            threading.Thread(target=_relay, daemon=True).start()
            return f"✅ Windows 通知攔截已啟動（{duration}s）" + (f"，過濾：{filter_app}" if filter_app else "")
    except Exception as e: return f"❌ 通知攔截失敗：{e}"


def execute_win_update(action):
    try:
        import subprocess
        if action in ("list", "check"):
            result = subprocess.run(["powershell", "-Command",
                "Get-WindowsUpdate -AcceptAll -Verbose 2>&1 | Select-Object -First 20"],
                capture_output=True, text=True, timeout=60)
            out = result.stdout.strip() or "無可用更新或需要 PSWindowsUpdate 模組"
            return f"🔄 Windows Update：\n{out[:1500]}"
        elif action == "install":
            result = subprocess.run(["powershell", "-Command",
                "Install-WindowsUpdate -AcceptAll -AutoReboot:$false -Verbose 2>&1"],
                capture_output=True, text=True, timeout=300)
            return f"✅ 更新執行完成：\n{result.stdout.strip()[:1500]}"
    except subprocess.TimeoutExpired:
        return "⏳ 更新查詢超時，請稍後再試"
    except Exception as e:
        return f"❌ Windows Update 失敗：{e}"


def execute_window_control(action, keyword=""):
    try:
        if action == "list":
            wins = [w for w in gw.getAllWindows() if w.title.strip()]
            return "\n".join(f"[{w._hWnd}] {w.title}" for w in wins) or "沒有視窗"
        wins = [w for w in gw.getAllWindows() if keyword.lower() in w.title.lower()]
        if not wins:
            return f"找不到視窗：{keyword}"
        w = wins[0]
        if action == "focus": w.activate(); return f"已切換到：{w.title}"
        elif action == "close": w.close(); return f"已關閉：{w.title}"
        elif action == "minimize": w.minimize(); return f"已最小化：{w.title}"
        elif action == "maximize": w.maximize(); return f"已最大化：{w.title}"
    except Exception as e:
        return f"執行失敗：{e}"


def execute_workflow(action, name="", steps=""):
    import json
    WORKFLOW_DIR = Path("C:/Users/blue_/workflows"); WORKFLOW_DIR.mkdir(exist_ok=True)
    if action == "list":
        files = list(WORKFLOW_DIR.glob("*.json"))
        return "\n".join(f.stem for f in files) if files else "沒有儲存的流程"
    elif action == "save":
        path = WORKFLOW_DIR / f"{name}.json"; data = json.loads(steps)
        path.write_text(json.dumps(data, ensure_ascii=False, indent=2), encoding="utf-8")
        return f"流程 [{name}] 已儲存"
    elif action == "run":
        path = WORKFLOW_DIR / f"{name}.json"
        if not path.exists(): return f"找不到流程：{name}"
        step_list = json.loads(path.read_text(encoding="utf-8"))
        for i, step in enumerate(step_list, 1):
            tool = step.get("tool"); args = step.get("args", []); delay = step.get("delay", 0)
            if delay: time.sleep(delay)
            try:
                tool_map = {"click": lambda a: pyautogui.click(int(a[0]), int(a[1])), "type": lambda a: pyautogui.write(" ".join(a), interval=0.05),
                    "press": lambda a: pyautogui.press(a[0]), "hotkey": lambda a: pyautogui.hotkey(*a),
                    "wait": lambda a: time.sleep(float(a[0])), "open": lambda a: subprocess.Popen(" ".join(a), shell=True)}
                if tool in tool_map: tool_map[tool](args)
            except Exception as e: return f"步驟 {i} 失敗：{e}"
        return f"流程 [{name}] 執行完畢（{len(step_list)} 步）"


def fetch_academic_search(query: str, field: str = "", lang: str = "en") -> str:
    try:
        from ddgs import DDGS
        region = "us-en" if lang == "en" else "tw-tzh"
        field_tag = f" {field}" if field else ""
        queries = [
            f"{query}{field_tag} site:scholar.google.com OR site:pubmed.ncbi.nlm.nih.gov OR site:semanticscholar.org",
            f"{query}{field_tag} research study findings",
            f"{query}{field_tag} academic paper review",
        ]
        all_hits = []
        with DDGS() as ddgs:
            for q in queries:
                hits = list(ddgs.text(q, region=region, max_results=4))
                all_hits.extend(hits)
        seen, lines = set(), [f"🎓 學術搜尋：{query}\n"]
        if field:
            lines.append(f"領域：{field}\n")
        count = 0
        for h in all_hits:
            title = h.get("title", "")
            if title in seen:
                continue
            seen.add(title)
            body = h.get("body", "")[:250]
            url = h.get("href", "")
            lines.append(f"📄 {title}\n   {body}")
            if url:
                lines.append(f"   {url}")
            lines.append("")
            count += 1
            if count >= 6:
                break
        lines.append(f"共找到 {count} 篇相關學術資料（建議至 Google Scholar 完整查閱）")
        return "\n".join(lines)
    except Exception as e:
        return f"學術搜尋失敗：{e}"


def fetch_analogy_maker(concept: str, audience: str = "一般大眾", count: int = 3, lang: str = "zh-tw") -> str:
    try:
        from ddgs import DDGS
        region = "tw-tzh" if lang == "zh-tw" else "us-en"
        with DDGS() as ddgs:
            hits = list(ddgs.text(f"{concept} 類比 比喻 解釋 說明", region=region, max_results=4))
        context = "\n".join(f"{h.get('title','')}：{h.get('body','')[:200]}" for h in hits[:4])
        lines = [
            f"🎯 類比說明：{concept}",
            f"目標受眾：{audience}\n",
            f"── 類比方案 ──\n",
        ]
        # 從搜尋結果找已有的類比，並提示生成
        existing = []
        for h in hits:
            body = h.get("body","")
            if any(kw in body for kw in ["就像", "好比", "類似", "如同", "比喻", "像是", "like", "similar to"]):
                existing.append(body[:200])
        for i, ex in enumerate(existing[:count], 1):
            lines.append(f"類比 {i}：{ex}\n")
        if len(existing) < count:
            lines.append(f"（以上為搜尋到的現有類比，Claude 會在回應中補充更多針對「{audience}」的類比說明）")
        lines.append(f"\n── 原始資料 ──")
        for h in hits[:2]:
            lines.append(f"• {h.get('title','')}：{h.get('body','')[:150]}")
        return "\n".join(lines)
    except Exception as e:
        return f"類比說明失敗：{e}"


def fetch_analyst_ratings(symbol: str) -> str:
    try:
        import yfinance as yf
        ticker = yf.Ticker(symbol)
        name = ticker.info.get("shortName") or symbol

        upgrades = ticker.upgrades_downgrades
        if upgrades is None or upgrades.empty:
            return f"{symbol} 無分析師評級資料"

        upgrades = upgrades.reset_index()
        recent = upgrades.head(10)
        lines = [f"📋 {name}（{symbol}）分析師評級\n"]
        for _, row in recent.iterrows():
            date = str(row.get("GradeDate", ""))[:10]
            firm = row.get("Firm", "")
            action = row.get("Action", "")
            to_grade = row.get("ToGrade", "")
            from_grade = row.get("FromGrade", "")
            emoji = "📈" if "up" in str(action).lower() or "initiat" in str(action).lower() else ("📉" if "down" in str(action).lower() else "➡️")
            lines.append(f"{emoji} {date} {firm}：{from_grade}→{to_grade}（{action}）")

        price = ticker.info.get("currentPrice", 0)
        target = ticker.info.get("targetMeanPrice", 0)
        if price and target:
            upside = (target / price - 1) * 100
            lines.append(f"\n分析師均價目標：{target:.2f}（現價 {price:.2f}，空間 {upside:+.1f}%）")
        return "\n".join(lines)
    except Exception as e:
        return f"分析師評級查詢失敗：{e}"


def fetch_app_navigator(app: str, task: str, input_text: str = "", monitor: int = 1, contact_name: str = "") -> str:
    try:
        import time
        import pyautogui
        import re
        app_lower = app.lower()
        results = []

        # 取得螢幕絕對座標偏移
        try:
            import mss
            with mss.mss() as sct:
                mons = sct.monitors
                mon_info = mons[monitor] if monitor < len(mons) else mons[1]
                mon_left = mon_info["left"]
                mon_top = mon_info["top"]
                mon_cx = mon_left + mon_info["width"] // 2
                mon_cy = mon_top + mon_info["height"] // 2
        except Exception:
            mon_left, mon_top, mon_cx, mon_cy = 0, 0, 960, 540

        # ══ Telegram 純螢幕控制（最優先，跳過所有 win32gui 視窗管理）══
        if "telegram" in app_lower:
            import pyperclip as _pc
            import ctypes as _ct, ctypes.wintypes as _wt
            import anthropic as _ant, base64 as _b64, io as _io2
            import ctypes as _ct2, win32ui as _w32u, win32con as _w32c2, win32gui as _w32g2
            import mss as _mss2
            from PIL import Image as _PI

            # 提取聯絡人名稱（優先用 contact_name 直接參數，再引號，再 regex）
            if contact_name:
                name = contact_name.strip()
            else:
                name_match = re.search(r'[「"](.*?)[」"]', task)
                if name_match:
                    name = name_match.group(1)
                else:
                    # 優先：「找/搜尋 XXX」格式，停於並/和/的/給/，/說/聊/打
                    m2 = re.search(r'(?:找到?|搜尋)(.+?)(?:並|和|的|給|，|,|說|聊|打|$)', task)
                    if m2:
                        name = m2.group(1).strip()
                    else:
                        # 次選：「跟 XXX 聊/說/打」格式，提取跟和動詞之間的名稱
                        m3 = re.search(r'跟(.+?)(?:聊|說|打|和|並|的|給|，|,|$)', task)
                        if m3:
                            name = m3.group(1).strip()
                        else:
                            name = re.sub(r'(螢幕\d|從Telegram|從telegram|Telegram|telegram|打開|找到|找|搜尋|對話|訊息|聊天|跟他說.*|跟.*?說.*|和|給|傳)', '', task).strip()

            # 名稱為空就直接報錯，不貼錯內容進搜尋框
            if not name:
                return "❌ 無法提取聯絡人名稱，請用引號標明，例如「去找「巴斯」聊天」或加 contact_name 參數"

            # ── 虛擬桌面參數（強制物理 DPI context，SendInput 用物理座標）──
            _u32 = _ct.windll.user32
            _old_ctx = _u32.SetThreadDpiAwarenessContext(_ct.c_void_p(-1))
            _vl = _u32.GetSystemMetrics(76); _vt = _u32.GetSystemMetrics(77)
            _vw = _u32.GetSystemMetrics(78); _vh = _u32.GetSystemMetrics(79)
            _u32.SetThreadDpiAwarenessContext(_old_ctx)
            # 取得螢幕邏輯尺寸，計算 DPI scale
            with _mss2.mss() as _smss:
                _mon_log = _smss.monitors[monitor]
                _log_left = _mon_log["left"]; _log_top = _mon_log["top"]
                _log_w = _mon_log["width"];   _log_h = _mon_log["height"]

            def _si_click(ax, ay):
                # ax, ay 是物理絕對座標（與 GetSystemMetrics DPI-aware 一致）
                class _MI(_ct.Structure):
                    _fields_ = [('dx',_wt.LONG),('dy',_wt.LONG),('mouseData',_wt.DWORD),
                                 ('dwFlags',_wt.DWORD),('time',_wt.DWORD),
                                 ('dwExtraInfo',_ct.POINTER(_ct.c_ulong))]
                class _U(_ct.Union):
                    _fields_ = [('mi',_MI)]
                class _INP(_ct.Structure):
                    _anonymous_ = ('u',); _fields_ = [('type',_wt.DWORD),('u',_U)]
                def _send(flags, dx=0, dy=0):
                    i = _INP(0,_U(mi=_MI(dx,dy,0,flags,0,None)))
                    _u32.SendInput(1,_ct.byref(i),_ct.sizeof(i))
                nx = int((ax - _vl) * 65535 // _vw)
                ny = int((ay - _vt) * 65535 // _vh)
                _send(0x0001|0x8000|0x4000, nx, ny)  # MOVE|ABSOLUTE|VIRTUALDESK
                time.sleep(0.35)
                _send(0x0002); time.sleep(0.08); _send(0x0004)  # DOWN + UP

            # ── GDI BitBlt 截圖（螢幕2用，其他用 dxcam）──
            def _cap(mon_num):
                with _mss2.mss() as s:
                    m = s.monitors[mon_num]
                    ml, mt, mw, mh = m["left"], m["top"], m["width"], m["height"]
                if mon_num == 2:
                    u32 = _ct2.windll.user32
                    old = u32.SetThreadDpiAwarenessContext(_ct2.c_void_p(-1))
                    try:
                        hd = _w32g2.GetDesktopWindow()
                        hdc = _w32g2.GetWindowDC(hd)
                        mdc = _w32u.CreateDCFromHandle(hdc)
                        sdc = mdc.CreateCompatibleDC()
                        bm = _w32u.CreateBitmap()
                        bm.CreateCompatibleBitmap(mdc, mw, mh)
                        sdc.SelectObject(bm)
                        sdc.BitBlt((0,0),(mw,mh),mdc,(ml,mt),_w32c2.SRCCOPY)
                        inf = bm.GetInfo(); bits = bm.GetBitmapBits(True)
                        img = _PI.frombuffer("RGB",(inf["bmWidth"],inf["bmHeight"]),bits,"raw","BGRX",0,1)
                        _w32g2.DeleteObject(bm.GetHandle())
                        sdc.DeleteDC(); mdc.DeleteDC(); _w32g2.ReleaseDC(hd, hdc)
                    finally:
                        u32.SetThreadDpiAwarenessContext(old)
                else:
                    import dxcam as _dx2
                    c = _dx2.create(output_idx={1:0,3:1}.get(mon_num,0))
                    img = _PI.fromarray(c.grab()); del c
                # 回傳 img + 邏輯偏移 + 邏輯寬（供計算 DPI scale）
                return img, ml, mt, mw

            # 判斷 GetSystemMetrics 是否物理座標
            _is_phys_gm = _vw > 4000

            # ── 截圖 → Claude Vision → 回傳與 GetSystemMetrics 同空間的絕對座標 ──
            def _see(prompt, mon_num):
                img, off_x_log, off_y_log, log_w = _cap(mon_num)
                # DPI scale：GDI 物理圖像寬 / mss 邏輯寬
                dpi = img.width / log_w  # 例：2400/1920=1.25
                ow = img.width
                if img.width > 1280:
                    r = 1280 / img.width
                    img = img.resize((1280, int(img.height*r)), _PI.LANCZOS)
                scale = ow / img.width  # 物理像素/resized像素
                buf = _io2.BytesIO()
                img.save(buf, format="JPEG", quality=85)
                b = _b64.standard_b64encode(buf.getvalue()).decode()
                resp = _ant.Anthropic().messages.create(
                    model="claude-haiku-4-5-20251001", max_tokens=80,
                    messages=[{"role":"user","content":[
                        {"type":"image","source":{"type":"base64","media_type":"image/jpeg","data":b}},
                        {"type":"text","text":f"截圖({img.width}x{img.height}px)。找「{prompt}」中心座標。僅回傳JSON:{{\"x\":整數,\"y\":整數,\"ok\":true/false}}"}
                    ]}]
                )
                import json as _j, re as _r2
                m2 = _r2.search(r'\{.*?\}', resp.content[0].text, _r2.DOTALL)
                if not m2: return None, None
                d = _j.loads(m2.group())
                if not d.get("ok", True): return None, None
                if _is_phys_gm:
                    # GetSystemMetrics=物理 → 回傳物理絕對
                    return round(off_x_log*dpi) + int(d["x"]*scale), \
                           round(off_y_log*dpi) + int(d["y"]*scale)
                else:
                    # GetSystemMetrics=邏輯 → pixel/dpi 轉邏輯再加邏輯偏移
                    return off_x_log + int(d["x"]*scale/dpi), \
                           off_y_log + int(d["y"]*scale/dpi)

            # ══ Step 0：先用 Vision 找 Telegram 視窗邊界（不管移到哪都能定位）══
            # Telegram 視覺特徵：左側深色聊天列表側欄、頂部≡選單+搜尋圖示、藍色紙飛機Logo
            def _find_tg_window(mon_num):
                """回傳 Telegram 視窗在螢幕上的絕對座標 (win_x, win_y, win_w, win_h)，失敗回傳 None"""
                img, off_x_log, off_y_log, log_w = _cap(mon_num)
                dpi = img.width / log_w
                ow = img.width
                if img.width > 1280:
                    r = 1280 / img.width
                    img = img.resize((1280, int(img.height*r)), _PI.LANCZOS)
                scale = ow / img.width
                buf = _io2.BytesIO(); img.save(buf, format="JPEG", quality=85)
                b = _b64.standard_b64encode(buf.getvalue()).decode()
                resp = _ant.Anthropic().messages.create(
                    model="claude-haiku-4-5-20251001", max_tokens=120,
                    messages=[{"role":"user","content":[
                        {"type":"image","source":{"type":"base64","media_type":"image/jpeg","data":b}},
                        {"type":"text","text":
                            f"截圖({img.width}x{img.height}px)。"
                            "找 Telegram 桌面應用程式視窗。"
                            "識別特徵：左側有深色聊天列表側欄（含聯絡人名稱和頭像）、頂部有≡漢堡選單圖示和搜尋放大鏡、藍色紙飛機Logo、右側有聊天對話區域。"
                            "回傳視窗邊界 JSON（像素座標）：{\"x\":左上角x, \"y\":左上角y, \"w\":視窗寬度, \"h\":視窗高度, \"ok\":true/false}"}
                    ]}]
                )
                m2 = re.search(r'\{.*?\}', resp.content[0].text, re.DOTALL)
                if not m2: return None
                import json as _j2
                d = _j2.loads(m2.group())
                if not d.get("ok", True): return None
                # 轉換到與 GetSystemMetrics 同空間的絕對座標
                wx = int(d["x"] * scale); wy = int(d["y"] * scale)
                ww = int(d["w"] * scale); wh = int(d["h"] * scale)
                if _is_phys_gm:
                    return round(off_x_log*dpi)+wx, round(off_y_log*dpi)+wy, ww, wh
                else:
                    return off_x_log+int(wx/dpi), off_y_log+int(wy/dpi), int(ww/dpi), int(wh/dpi)

            _tg_win = _find_tg_window(monitor)
            if _tg_win:
                _tw_x, _tw_y, _tw_w, _tw_h = _tg_win
                results.append(f"📐 Telegram視窗({_tw_x},{_tw_y}) {_tw_w}x{_tw_h}")
            else:
                # 找不到視窗邊界時，退回用 monitor 左上角估算
                _tw_x = round(_log_left * 1.25) if _is_phys_gm else _log_left
                _tw_y = round(_log_top * 1.25) if _is_phys_gm else _log_top
                _tw_w = round(_log_w * 1.25) if _is_phys_gm else _log_w
                _tw_h = round(_log_h * 1.25) if _is_phys_gm else _log_h
                results.append(f"⚠️ 找不到視窗邊界，用螢幕偏移估算({_tw_x},{_tw_y})")

            # ── Step 0.5：確保 Telegram 在正常聊天列表狀態 ──
            import pyautogui as _pg
            _pg.press("escape"); time.sleep(0.3)
            _pg.press("escape"); time.sleep(0.3)
            # 檢查是否在異常狀態（設定面板、空白畫面等）
            try:
                import pytesseract as _tess_st
                _state_img, _, _ = _cap(monitor)
                _state_text = _tess_st.image_to_string(_state_img, lang="chi_tra+eng")
                if "請選擇聊天對象" in _state_text or "我的資料" in _state_text or "設定" in _state_text or "建立群組" in _state_text:
                    _pg.press("escape"); time.sleep(0.3)
                    _pg.press("escape"); time.sleep(0.3)
                    results.append("🔄 從異常狀態恢復")
            except Exception:
                pass

            # ① 用鍵盤快捷鍵開搜尋（不用 Vision 找搜尋框，避免點到漢堡選單）
            _si_click(_tw_x + int(_tw_w * 0.13), _tw_y + int(_tw_h * 0.05))
            time.sleep(0.3)
            results.append("🔍 點擊搜尋區域")

            # ② 清空 + 貼上聯絡人名稱
            _pg.hotkey("ctrl","a"); time.sleep(0.1); _pg.press("delete"); time.sleep(0.1)
            _pc.copy(name); _pg.hotkey("ctrl","v"); time.sleep(1.2)
            results.append(f"🔍 搜尋名稱：{name}")

            # ③ 用 Vision 在搜尋結果中找聯絡人
            _rx, _ry = _see(
                f"Telegram搜尋結果列表中名稱為「{name}」的聊天項目（個人對話優先），點名字或頭像位置",
                monitor
            )
            if _rx is not None:
                _si_click(_rx, _ry); time.sleep(0.7)
                results.append(f"✅ Vision找到並點擊「{name}」({_rx},{_ry})")
            else:
                # 備用：搜尋框下方固定偏移點第一個結果
                _first_y = _tw_y + int(_tw_h * 0.11)
                _si_click(_tw_x + int(_tw_w * 0.13), _first_y); time.sleep(0.7)
                results.append(f"⚠️ Vision找不到，點第一個搜尋結果")

            # ④ 有訊息就輸入送出
            if input_text:
                time.sleep(0.5)
                mx, my = _see(
                    "Telegram聊天視窗底部的訊息輸入框（最下方打字區，顯示「輸入訊息」或有emoji圖示和迴紋針圖示，不是頂部搜尋框）",
                    monitor
                )
                if mx is None:
                    mx = _tw_x + int(_tw_w * 0.6); my = _tw_y + int(_tw_h * 0.95)
                    results.append(f"⚠️ Vision找不到訊息框，視窗相對備用({mx},{my})")
                else:
                    _si_click(mx, my); time.sleep(0.4)

                # 偵測編號清單（1. xxx\n2. xxx），拆成多則分開發送
                _lines = re.split(r'\n(?=\d+[\.\、\)])', input_text.strip())
                if len(_lines) > 1:
                    for _li, _line in enumerate(_lines):
                        _msg = re.sub(r'^\d+[\.\、\)]\s*', '', _line).strip()
                        if not _msg:
                            continue
                        _pc.copy(_msg); _pg.hotkey("ctrl","v"); time.sleep(0.2)
                        _pg.press("enter"); time.sleep(0.4)
                        results.append(f"📤 第{_li+1}則：{_msg[:20]}")
                else:
                    _pc.copy(input_text); _pg.hotkey("ctrl","v"); time.sleep(0.3)
                    _pg.press("enter"); time.sleep(0.2); _pg.press("enter")
                    results.append(f"📤 已送：{input_text}")

            # ── 清理：點擊聊天列表第一個聊天回到正常狀態（不用 Esc 避免關掉聊天室）──
            time.sleep(0.3)
            _si_click(_tw_x + int(_tw_w * 0.13), _tw_y + int(_tw_h * 0.12))
            time.sleep(0.2)
            results.append("🧹 已回到聊天列表")

            return "\n".join(results) if results else "Telegram導航完成"

        # ── 非 Telegram：通用視窗管理 ──
        # 先點一下目標螢幕中央，確保焦點在正確螢幕
        pyautogui.click(mon_cx, mon_cy)
        time.sleep(0.3)

        # 通用：把視窗拉到前景
        try:
            import win32gui, win32con
            def find_window(name):
                result = []
                win32gui.EnumWindows(
                    lambda h, _: result.append(h)
                    if name.lower() in win32gui.GetWindowText(h).lower() else None, None)
                return result[0] if result else None
            hw = find_window(app)
            if hw:
                win32gui.ShowWindow(hw, win32con.SW_RESTORE)
                win32gui.SetForegroundWindow(hw)
                time.sleep(0.6)
                results.append(f"✅ 切換到 {app} 視窗")
            else:
                results.append(f"⚠️ 找不到 {app} 視窗，嘗試開啟")
                pyautogui.hotkey("win", "s")
                time.sleep(0.5)
                pyautogui.write(app, interval=0.05)
                time.sleep(1)
                pyautogui.press("enter")
                time.sleep(2.5)
        except Exception as e:
            results.append(f"視窗切換：{e}")

        # 再次點目標螢幕確保焦點
        pyautogui.click(mon_cx, mon_cy)
        time.sleep(0.3)

        # LINE 專屬流程
        if "line" in app_lower:
            name_match = re.search(r'[「"](.*?)[」"]', task)
            name = name_match.group(1) if name_match else re.sub(r'(打開|找到|找|搜尋|對話|訊息)', '', task).strip()
            r = fetch_ocr_click(name, monitor)
            results.append(f"找對話 {name}：{r}")
            if input_text:
                time.sleep(0.5)
                pyautogui.write(input_text, interval=0.04)
                if "送出" in task or "回覆" in task:
                    pyautogui.press("enter")
                    results.append("📤 已送出")

        # 通用流程：用 OCR/視覺找目標
        else:
            r = fetch_ocr_click(task, monitor)
            results.append(f"OCR操作：{r}")
            if input_text:
                time.sleep(0.3)
                pyautogui.write(input_text, interval=0.04)
                results.append(f"已輸入：{input_text[:40]}")

        return "\n".join(results) if results else f"App導航完成：{app} / {task}"
    except Exception as e:
        return f"App導航失敗：{e}"


def fetch_ashare(code: str, period: str = "1mo") -> str:
    """A股（滬深）/ 港股查詢，自動判斷市場並加後綴"""
    try:
        import yfinance as yf
        code = code.strip().lstrip("0" * 0)  # 保留原始代號

        # 判斷市場：6位數字 → A股（6開頭=上海.SS，其他=深圳.SZ）；4位以下 → 港股.HK
        if code.isdigit():
            if len(code) == 6:
                symbol = f"{code}.SS" if code.startswith("6") else f"{code}.SZ"
                market = "A股"
            else:
                symbol = f"{code.zfill(4)}.HK"
                market = "港股"
        else:
            symbol = code  # 使用者直接提供後綴
            market = "未知"

        ticker = yf.Ticker(symbol)
        info = ticker.info
        hist = ticker.history(period="3mo")

        if hist.empty:
            return f"找不到「{code}」的數據，請確認代號是否正確。"

        name = info.get("longName") or info.get("shortName") or code
        currency = info.get("currency", "")
        current = hist["Close"].iloc[-1]
        prev = hist["Close"].iloc[-2] if len(hist) > 1 else current
        change = current - prev
        change_pct = (change / prev * 100) if prev else 0
        arrow = "▲" if change >= 0 else "▼"
        volume = hist["Volume"].iloc[-1]

        ma5  = hist["Close"].tail(5).mean()
        ma20 = hist["Close"].tail(20).mean()
        ma60 = hist["Close"].tail(60).mean() if len(hist) >= 60 else hist["Close"].mean()
        rsi  = calc_rsi(hist["Close"]) if len(hist) >= 15 else None

        if ma5 > ma20 > ma60:
            trend = "強勢多頭 📈"
        elif ma5 < ma20 < ma60:
            trend = "強勢空頭 📉"
        elif ma5 > ma20:
            trend = "短線偏多 🔼"
        else:
            trend = "短線偏空 🔽"

        rsi_note = ""
        if rsi is not None:
            if rsi >= 80:   rsi_note = "（嚴重超買 ⚠️）"
            elif rsi >= 70: rsi_note = "（超買）"
            elif rsi <= 20: rsi_note = "（嚴重超賣 💡）"
            elif rsi <= 30: rsi_note = "（超賣）"
            else:           rsi_note = "（中性）"

        week52_high = info.get("fiftyTwoWeekHigh")
        week52_low  = info.get("fiftyTwoWeekLow")
        market_cap  = info.get("marketCap")
        pe_ratio    = info.get("trailingPE")

        result = (
            f"🇨🇳 {name}（{market} {code}）\n"
            f"💰 現價：{current:.2f} {currency}  {arrow} {abs(change):.2f} ({change_pct:+.2f}%)\n"
            f"📦 成交量：{volume:,}\n"
            f"\n── 技術指標 ──\n"
            f"MA5：{ma5:.2f}　MA20：{ma20:.2f}　MA60：{ma60:.2f}\n"
            f"趨勢：{trend}\n"
        )
        if rsi is not None:
            result += f"RSI(14)：{rsi}{rsi_note}\n"
        if week52_high and week52_low:
            result += f"52週高低：{week52_low:.2f} ~ {week52_high:.2f}\n"
        result += "\n── 基本面 ──\n"
        if market_cap:
            mc_str = f"{market_cap/1e12:.2f}兆" if market_cap >= 1e12 else f"{market_cap/1e8:.0f}億"
            result += f"市值：{mc_str} {currency}\n"
        if pe_ratio:
            result += f"本益比：{pe_ratio:.1f}\n"

        return result.strip()
    except Exception as e:
        return f"查詢「{code}」失敗：{e}"


def fetch_asset_allocation(age: int, risk_level: str, goal: str = "退休",
                            investment_horizon: int = None) -> str:
    try:
        horizon = investment_horizon or max(65 - age, 5)
        # 基本股債比（110法則）
        base_stock = min(110 - age, 90)
        if risk_level == "保守":
            stock = max(base_stock - 20, 10)
        elif risk_level == "積極":
            stock = min(base_stock + 15, 90)
        else:
            stock = base_stock
        bond = max(100 - stock - 10, 0)
        cash = 100 - stock - bond

        # 細分配置
        tw_stock = round(stock * 0.4)
        us_stock = round(stock * 0.35)
        intl_stock = stock - tw_stock - us_stock
        tw_bond = round(bond * 0.3)
        us_bond = bond - tw_bond

        lines = [
            f"📊 資產配置建議\n",
            f"年齡：{age} 歲　風險偏好：{risk_level}",
            f"目標：{goal}　投資期間：{horizon} 年\n",
            f"── 大類配置 ──",
            f"股票：{stock}%　債券：{bond}%　現金：{cash}%\n",
            f"── 細分建議 ──",
            f"台股（0050/0056）：{tw_stock}%",
            f"美股（VTI/VOO）：{us_stock}%",
            f"國際股（VEA/VWO）：{intl_stock}%",
            f"台灣公債/ETF：{tw_bond}%",
            f"美債（BND/TLT）：{us_bond}%",
            f"現金/定存：{cash}%\n",
            f"── 再平衡建議 ──",
            f"每年或偏離5%以上時再平衡",
            f"隨年齡增長逐步降低股票比例",
        ]
        if risk_level == "保守":
            lines.append("\n⚠️ 保守型：優先保本，適合距退休較近者")
        elif risk_level == "積極":
            lines.append("\n💡 積極型：承受較大波動換取長期成長，需有10年以上視野")
        return "\n".join(lines)
    except Exception as e:
        return f"資產配置建議失敗：{e}"


def fetch_backtest(symbol: str, strategy: str = "ma_cross", period: str = "2y") -> str:
    """回測投資策略"""
    try:
        import yfinance as yf
        import numpy as np

        ticker = yf.Ticker(symbol)
        hist = ticker.history(period=period)
        if hist.empty or len(hist) < 30:
            return f"找不到 {symbol} 的歷史資料"

        name = ticker.info.get("shortName") or symbol
        close = hist["Close"]
        n = len(close)

        if strategy == "buy_hold":
            total_ret = (close.iloc[-1] / close.iloc[0] - 1) * 100
            annual_ret = ((close.iloc[-1] / close.iloc[0]) ** (252 / n) - 1) * 100
            max_dd = ((close / close.cummax()) - 1).min() * 100
            return (
                f"📈 {name}（{symbol}）買進持有回測（{period}）\n\n"
                f"起始價：{close.iloc[0]:.2f}\n"
                f"結束價：{close.iloc[-1]:.2f}\n"
                f"總報酬：{total_ret:+.2f}%\n"
                f"年化報酬：{annual_ret:+.2f}%\n"
                f"最大回撤：{max_dd:.2f}%"
            )

        elif strategy == "ma_cross":
            ma5  = close.rolling(5).mean()
            ma20 = close.rolling(20).mean()
            signal = (ma5 > ma20).astype(int)
            signal_prev = signal.shift(1)
            buy_signals  = (signal == 1) & (signal_prev == 0)
            sell_signals = (signal == 0) & (signal_prev == 1)

            trades = []
            buy_price = None
            for i in range(len(close)):
                if buy_signals.iloc[i] and buy_price is None:
                    buy_price = close.iloc[i]
                elif sell_signals.iloc[i] and buy_price is not None:
                    ret = (close.iloc[i] / buy_price - 1) * 100
                    trades.append(ret)
                    buy_price = None

            if not trades:
                return f"{symbol} 在 {period} 內無均線交叉訊號"

            wins = sum(1 for t in trades if t > 0)
            total_ret = sum(trades)
            avg_ret   = total_ret / len(trades)
            win_rate  = wins / len(trades) * 100
            max_dd    = ((close / close.cummax()) - 1).min() * 100

            return (
                f"📊 {name}（{symbol}）MA5穿MA20 策略回測（{period}）\n\n"
                f"交易次數：{len(trades)} 次\n"
                f"勝率：{win_rate:.1f}%\n"
                f"平均每筆報酬：{avg_ret:+.2f}%\n"
                f"累計報酬：{total_ret:+.2f}%\n"
                f"最大回撤：{max_dd:.2f}%\n\n"
                f"買進持有同期：{(close.iloc[-1]/close.iloc[0]-1)*100:+.2f}%"
            )

        elif strategy == "dca":
            # 每月第一個交易日買進固定金額
            monthly = close.resample("MS").first()
            shares = 0
            cost   = 0
            invest_per_month = 10000  # 每月投入 10000 元
            for price in monthly:
                shares += invest_per_month / price
                cost   += invest_per_month
            final_value = shares * close.iloc[-1]
            total_ret   = (final_value / cost - 1) * 100

            return (
                f"📅 {name}（{symbol}）定期定額回測（{period}，每月10000元）\n\n"
                f"總投入：{cost:,.0f} 元\n"
                f"最終市值：{final_value:,.0f} 元\n"
                f"總報酬：{total_ret:+.2f}%\n"
                f"累積股數：{shares:.2f} 股\n"
                f"平均成本：{cost/shares:.2f} 元"
            )

        return "不支援的策略"
    except Exception as e:
        return f"回測失敗：{e}"


def fetch_benchmark_analysis(subject: str, industry: str = "", lang: str = "zh-tw") -> str:
    try:
        from ddgs import DDGS
        region = "tw-tzh" if lang == "zh-tw" else "us-en"
        ind_tag = f" {industry}" if industry else ""
        queries = {
            "業界標竿": f"{industry or subject} 最佳實踐 業界標準 領導者",
            "對標案例": f"{subject}{ind_tag} 對標 學習 比較 領先",
            "改進方向": f"{subject}{ind_tag} 改善 優化 提升 差距",
        }
        lines = [f"📐 標竿分析：{subject}"]
        if industry:
            lines.append(f"產業：{industry}")
        lines.append("")
        with DDGS() as ddgs:
            for sec, q in queries.items():
                hits = list(ddgs.text(q, region=region, max_results=3))
                lines.append(f"── {sec} ──")
                for h in hits[:2]:
                    lines.append(f"• {h.get('title','')}：{h.get('body','')[:200]}")
                lines.append("")
        lines.append("── 建議行動 ──")
        lines.append(f"1. 找出領域內 Top 3 標竿對象，深入研究其核心做法")
        lines.append(f"2. 識別 {subject} 與標竿的具體差距")
        lines.append(f"3. 制定可量化的追趕目標與時間表")
        return "\n".join(lines)
    except Exception as e:
        return f"標竿分析失敗：{e}"


def fetch_bias_detector(topic: str, lang: str = "zh-tw") -> str:
    try:
        from ddgs import DDGS
        region = "tw-tzh" if lang == "zh-tw" else "us-en"
        queries = [
            f"{topic} 媒體偏見 立場",
            f"{topic} 支持方 反對方 爭議",
            f"{topic} 批評 質疑 偏頗",
        ]
        all_hits = []
        with DDGS() as ddgs:
            for q in queries:
                hits = list(ddgs.text(q, region=region, max_results=3))
                all_hits.extend(hits)
        # 分析立場傾向
        pro_kw = ["支持", "推崇", "認同", "正面", "優秀", "progressive", "support", "advocate"]
        con_kw = ["反對", "批評", "質疑", "保守", "警告", "oppose", "criticize", "concern"]
        left_kw = ["進步", "左派", "liberal", "progressive", "left"]
        right_kw = ["保守", "右派", "conservative", "right", "traditional"]
        all_text = " ".join(h.get("body","") + h.get("title","") for h in all_hits).lower()
        pro_s = sum(all_text.count(k) for k in pro_kw)
        con_s = sum(all_text.count(k) for k in con_kw)
        left_s = sum(all_text.count(k) for k in left_kw)
        right_s = sum(all_text.count(k) for k in right_kw)
        bias_dir = "中立" if abs(pro_s - con_s) < 3 else ("偏正面/支持" if pro_s > con_s else "偏負面/批評")
        political = "中立" if abs(left_s - right_s) < 2 else ("偏進步/左派" if left_s > right_s else "偏保守/右派")
        lines = [
            f"🔍 偏見偵測：{topic}\n",
            f"── 立場分析 ──",
            f"情感傾向：{bias_dir}（正面訊號{pro_s} vs 負面訊號{con_s}）",
            f"政治傾向：{political}",
            f"",
            f"── 各方觀點樣本 ──",
        ]
        seen = set()
        for h in all_hits[:6]:
            title = h.get("title","")
            if title in seen:
                continue
            seen.add(title)
            body = h.get("body","")[:150]
            lines.append(f"• {title}：{body}")
        lines.append(f"\n⚠️ 閱讀此議題資料時建議多方比對，避免單一來源形成片面認知")
        return "\n".join(lines)
    except Exception as e:
        return f"偏見偵測失敗：{e}"


def fetch_bond_yield() -> str:
    """美國公債殖利率"""
    try:
        import yfinance as yf
        bonds = {
            "3個月": "^IRX",
            "2年":   "2YY=F",
            "5年":   "^FVX",
            "10年":  "^TNX",
            "30年":  "^TYX",
        }
        lines = ["🏦 美國公債殖利率\n"]
        yields = {}
        for label, sym in bonds.items():
            try:
                hist = yf.Ticker(sym).history(period="2d")
                if hist.empty:
                    continue
                val = hist["Close"].iloc[-1]
                yields[label] = val
                if len(hist) >= 2:
                    chg = val - hist["Close"].iloc[-2]
                    arrow = "▲" if chg >= 0 else "▼"
                    lines.append(f"{label}：{val:.3f}%  {arrow} {chg:+.3f}%")
                else:
                    lines.append(f"{label}：{val:.3f}%")
            except Exception:
                pass

        # 利差分析
        if "2年" in yields and "10年" in yields:
            spread = yields["10年"] - yields["2年"]
            curve = "正斜率（景氣正常）" if spread > 0 else "倒掛（衰退警訊 ⚠️）"
            lines.append(f"\n10Y-2Y 利差：{spread:+.3f}%  {curve}")
        if "3個月" in yields and "10年" in yields:
            spread2 = yields["10年"] - yields["3個月"]
            curve2 = "正常" if spread2 > 0 else "倒掛 ⚠️"
            lines.append(f"10Y-3M 利差：{spread2:+.3f}%  {curve2}")

        return "\n".join(lines)
    except Exception as e:
        return f"殖利率查詢失敗：{e}"


def fetch_brainstorm(problem: str, count: int = 8, style: str = "實用", lang: str = "zh-tw") -> str:
    try:
        from ddgs import DDGS
        region = "tw-tzh" if lang == "zh-tw" else "us-en"
        style_map = {
            "實用": "解決方案 方法 做法",
            "創意": "創新 新穎 非傳統 獨特",
            "顛覆": "顛覆 革命性 完全不同 打破慣例",
        }
        style_q = style_map.get(style, "解決方案 方法")
        all_hits = []
        with DDGS() as ddgs:
            for q in [f"{problem} {style_q}", f"{problem} 案例 成功", f"{problem} 創意 想法"]:
                hits = list(ddgs.text(q, region=region, max_results=3))
                all_hits.extend(hits)
        lines = [f"🧠 腦力激盪：{problem}", f"風格：{style}　目標：{count} 個方案\n"]
        # 從搜尋結果萃取方向
        seen, ideas = set(), []
        for h in all_hits:
            title = h.get("title","")
            body = h.get("body","")
            if title not in seen and body:
                seen.add(title)
                ideas.append(f"{title}：{body[:120]}")
        lines.append("── 方案清單 ──")
        for i, idea in enumerate(ideas[:count], 1):
            lines.append(f"{i}. {idea}\n")
        if len(ideas) < count:
            lines.append(f"（已蒐集 {len(ideas)} 個方向，可進一步細化）")
        return "\n".join(lines)
    except Exception as e:
        return f"腦力激盪失敗：{e}"


def fetch_china_search(query: str, category: str = "其他", count: int = 6) -> str:
    """全方位中國大陸資訊搜尋：旅遊/美食/文化/戲劇/演員/工作等"""
    try:
        count = min(max(count, 1), 10)
        results = []

        # 1. Google News（中文簡體，抓最新新聞/資訊）
        try:
            import feedparser
            news_query = query
            # 依分類補充關鍵字讓搜尋更精準
            category_hints = {
                "旅遊": f"{query} 旅遊攻略 景點",
                "美食": f"{query} 美食 餐廳 推薦",
                "文化風俗": f"{query} 文化 習俗 傳統",
                "戲劇影視": f"{query} 電視劇 電影 劇情",
                "演員明星": f"{query} 演員 明星 近況",
                "工作生活": f"{query} 工作 薪資 生活",
                "城市介紹": f"{query} 城市 介紹 特色",
                "歷史": f"{query} 歷史 背景",
                "科技": f"{query} 科技 技術",
                "新聞時事": f"{query} 最新 新聞",
            }
            news_query = category_hints.get(category, query)
            url = f"https://news.google.com/rss/search?q={urllib.parse.quote(news_query)}&hl=zh-Hans&gl=CN&ceid=CN:zh-Hans"
            feed = feedparser.parse(url)
            if feed.entries:
                lines = [f"📰 Google 新聞（{category}）"]
                for i, entry in enumerate(feed.entries[:min(count, 5)], 1):
                    title = entry.get("title", "").split(" - ")[0]  # 去掉媒體名稱
                    pub = entry.get("published", "")[:16]
                    lines.append(f"{i}. {title}（{pub}）")
                results.append("\n".join(lines))
        except Exception:
            pass

        # 2. DuckDuckGo 搜尋（zh-cn 地區，含陸網資料）
        try:
            from ddgs import DDGS
            ddg_results = []
            with DDGS() as ddgs:
                for r in ddgs.text(query, region="zh-cn", max_results=count):
                    title = r.get("title", "")
                    body  = r.get("body", "")[:150]
                    href  = r.get("href", "")
                    ddg_results.append(f"• {title}\n  {body}\n  {href}")
            if ddg_results:
                results.append(f"🔍 網路搜尋結果\n" + "\n\n".join(ddg_results))
        except Exception:
            pass

        # 3. Wikipedia 中文（適合文化/歷史/人物類）
        if category in ("文化風俗", "歷史", "演員明星", "城市介紹", "戲劇影視"):
            try:
                wiki_url = f"https://zh.wikipedia.org/api/rest_v1/page/summary/{urllib.parse.quote(query.split()[0])}"
                resp = requests.get(wiki_url, timeout=8)
                if resp.status_code == 200:
                    data = resp.json()
                    extract = data.get("extract", "")
                    title_w = data.get("title", "")
                    if extract:
                        results.append(f"📖 Wikipedia：{title_w}\n{extract[:500]}")
            except Exception:
                pass

        if not results:
            return f"找不到「{query}」的相關資訊，請嘗試更換關鍵字"

        return "\n\n" + "─" * 25 + "\n\n".join(results)

    except Exception as e:
        return f"中國資訊搜尋失敗：{e}"


def fetch_cn_news(source: str = "all", count: int = 5) -> str:
    """抓取中國大陸新聞 RSS"""
    try:
        import feedparser
        count = min(max(count, 1), 10)
        feeds = {
            "xinhua":  ("新華社",   "https://feeds.xinhuanet.com/news/world/rss"),
            "people":  ("人民網",   "http://www.people.com.cn/rss/politics.xml"),
            "36kr":    ("36氪",     "https://36kr.com/feed"),
            "caixin":  ("財新網",   "https://www.caixin.com/rss/home.xml"),
        }
        # 備用可靠來源
        fallback_feeds = {
            "xinhua":  ("新華社",   "https://rsshub.app/xinhua/world"),
            "people":  ("人民網",   "https://rsshub.app/people/politics"),
            "36kr":    ("36氪",     "https://rsshub.app/36kr/news/latest"),
            "caixin":  ("財新網",   "https://rsshub.app/caixin/blog"),
        }
        sources = list(feeds.keys()) if source == "all" else [source]
        results = []
        for src in sources:
            if src not in feeds:
                continue
            label, url = feeds[src]
            try:
                feed = feedparser.parse(url)
                items = feed.entries[:count]
                if not items:
                    # 嘗試備用
                    label2, url2 = fallback_feeds.get(src, (label, url))
                    feed = feedparser.parse(url2)
                    items = feed.entries[:count]
                if not items:
                    results.append(f"📰 {label}：暫無資料")
                    continue
                lines = [f"📰 {label}"]
                for i, entry in enumerate(items, 1):
                    title = entry.get("title", "無標題")
                    lines.append(f"{i}. {title}")
                results.append("\n".join(lines))
            except Exception:
                results.append(f"📰 {label}：抓取失敗")
        return "\n\n".join(results) if results else "無法取得中國新聞"
    except Exception as e:
        return f"中國新聞查詢失敗：{e}"


def fetch_commodity(items: list = None) -> str:
    """黃金/原油/原物料報價"""
    try:
        import yfinance as yf
        commodity_map = {
            "gold":   ("黃金",   "GC=F",  "USD/盎司"),
            "oil":    ("WTI原油", "CL=F",  "USD/桶"),
            "silver": ("白銀",   "SI=F",  "USD/盎司"),
            "copper": ("銅",     "HG=F",  "USD/磅"),
            "natgas": ("天然氣", "NG=F",  "USD/MMBtu"),
            "wheat":  ("小麥",   "ZW=F",  "USd/英斗"),
            "corn":   ("玉米",   "ZC=F",  "USd/英斗"),
        }
        if not items or "all" in items:
            items = list(commodity_map.keys())

        lines = ["🛢 大宗商品報價\n"]
        for key in items:
            if key not in commodity_map:
                continue
            name, sym, unit = commodity_map[key]
            try:
                hist = yf.Ticker(sym).history(period="2d")
                if hist.empty:
                    continue
                price = hist["Close"].iloc[-1]
                if len(hist) >= 2:
                    chg = (price / hist["Close"].iloc[-2] - 1) * 100
                    arrow = "▲" if chg >= 0 else "▼"
                    lines.append(f"{name}：{price:.2f} {unit}  {arrow} {chg:+.2f}%")
                else:
                    lines.append(f"{name}：{price:.2f} {unit}")
            except Exception:
                pass
        return "\n".join(lines)
    except Exception as e:
        return f"商品報價失敗：{e}"


def fetch_company_research(company: str, lang: str = "zh-tw") -> str:
    try:
        from ddgs import DDGS
        import yfinance as yf
        region = "tw-tzh" if lang == "zh-tw" else "us-en"
        lines = [f"🏢 公司研究：{company}\n"]
        # 嘗試取得財務數據
        try:
            ticker = yf.Ticker(company)
            info = ticker.info
            name = info.get("longName") or info.get("shortName") or company
            sector = info.get("sector") or "—"
            industry = info.get("industry") or "—"
            employees = info.get("fullTimeEmployees") or "—"
            revenue = info.get("totalRevenue")
            market_cap = info.get("marketCap")
            lines += [
                f"── 基本資料 ──",
                f"公司：{name}　產業：{sector} / {industry}",
            ]
            if employees != "—":
                lines.append(f"員工數：{employees:,}")
            if revenue:
                lines.append(f"年營收：${revenue/1e8:.1f}億")
            if market_cap:
                lines.append(f"市值：${market_cap/1e8:.1f}億")
            lines.append("")
        except Exception:
            pass
        # 搜尋新聞與評價
        queries = {
            "公司動態": f"{company} 最新消息 發展",
            "產品評價": f"{company} 產品 服務 評價",
            "競爭分析": f"{company} 競爭對手 市場地位",
        }
        with DDGS() as ddgs:
            for sec, q in queries.items():
                hits = list(ddgs.text(q, region=region, max_results=2))
                if hits:
                    lines.append(f"── {sec} ──")
                    for h in hits:
                        lines.append(f"• {h.get('title','')}：{h.get('body','')[:180]}")
                    lines.append("")
        return "\n".join(lines)
    except Exception as e:
        return f"公司研究失敗：{e}"


def fetch_compare_analysis(items: list, dimensions: list = None, context: str = "") -> str:
    try:
        from ddgs import DDGS
        ctx = f" {context}" if context else ""
        data = {}
        with DDGS() as ddgs:
            for item in items[:5]:
                hits = list(ddgs.text(f"{item}{ctx} 評價 特點 優缺點", region="tw-tzh", max_results=3))
                data[item] = " ".join(h.get("body","") for h in hits)[:500]
        lines = [f"⚖️ 比較分析：{' vs '.join(items)}\n"]
        if context:
            lines.append(f"背景：{context}\n")
        for item in items[:5]:
            lines.append(f"【{item}】")
            lines.append(data.get(item, "資料不足")[:400])
            lines.append("")
        lines.append("── 綜合建議 ──")
        lines.append(f"以上為各項資料彙整，請根據您的需求與優先考量做最終選擇。")
        return "\n".join(lines)
    except Exception as e:
        return f"比較分析失敗：{e}"


def fetch_compound_calculator(principal: float, annual_rate: float, years: int,
                               monthly_add: float = 0, compound_freq: int = 12) -> str:
    try:
        compound_freq = int(compound_freq) if not isinstance(compound_freq, int) else compound_freq
        if compound_freq <= 0:
            compound_freq = 12
        r = annual_rate / 100 / compound_freq
        n = compound_freq * years
        # 本金複利
        future_principal = principal * (1 + r) ** n
        # 每月定期投入（按月計算）
        if monthly_add > 0:
            r_m = annual_rate / 100 / 12
            n_m = years * 12
            future_monthly = monthly_add * (((1 + r_m) ** n_m - 1) / r_m) if r_m > 0 else monthly_add * n_m
        else:
            future_monthly = 0
        total = future_principal + future_monthly
        total_invest = principal + monthly_add * years * 12
        profit = total - total_invest
        lines = [
            f"📈 複利計算器\n",
            f"本金：{principal:,.0f} 元",
            f"年化報酬：{annual_rate}%　期間：{years} 年",
            f"每月追加：{monthly_add:,.0f} 元",
            f"",
            f"── 試算結果 ──",
            f"本金複利成長：{future_principal:,.0f} 元",
        ]
        if monthly_add > 0:
            lines.append(f"追加投入成長：{future_monthly:,.0f} 元")
        lines += [
            f"期末總資產：{total:,.0f} 元",
            f"總投入成本：{total_invest:,.0f} 元",
            f"獲利：{profit:,.0f} 元（{profit/total_invest*100:.1f}%）",
            f"",
            f"── 72法則 ──",
            f"資產翻倍需：{72/annual_rate:.1f} 年（年報酬 {annual_rate}%）",
        ]
        return "\n".join(lines)
    except Exception as e:
        return f"複利計算失敗：{e}"


def fetch_concept_stocks(theme: str) -> str:
    try:
        # 台股概念股資料庫
        concepts = {
            "AI": ["2330.TW","2303.TW","2454.TW","3711.TW","2379.TW","2308.TW","6488.TW","3017.TW","2382.TW","5274.TW"],
            "電動車": ["2308.TW","1590.TW","2207.TW","6239.TW","1802.TW","2049.TW","3665.TW","1537.TW","6227.TW","2371.TW"],
            "軍工": ["1323.TW","2409.TW","2348.TW","8112.TW","1536.TW","2634.TW","6245.TW","1513.TW"],
            "低軌衛星": ["3045.TW","2230.TW","6438.TW","3413.TW","2365.TW","3508.TW","6285.TW","4306.TW"],
            "半導體": ["2330.TW","2303.TW","2454.TW","2308.TW","3711.TW","5347.TW","6770.TW","3034.TW","2379.TW","4919.TW"],
            "5G": ["2412.TW","3045.TW","2498.TW","2356.TW","3231.TW","6488.TW","2439.TW","3293.TW"],
            "儲能": ["1907.TW","1504.TW","6409.TW","3481.TW","2023.TW","6121.TW","1590.TW"],
            "DRAM": ["2303.TW","3450.TW","4967.TW","3260.TW","2408.TW"],
            "CoWoS": ["2330.TW","6235.TW","3711.TW","2454.TW","3036.TW","8046.TW","6415.TW"],
            "矽光子": ["2330.TW","3008.TW","2454.TW","6510.TW","3081.TW"],
            "機器人": ["2308.TW","1590.TW","2049.TW","1537.TW","3665.TW","6288.TW","2382.TW"],
            "航運": ["2603.TW","2609.TW","2615.TW","2610.TW","2618.TW","5608.TW"],
            "金融": ["2881.TW","2882.TW","2886.TW","2891.TW","2884.TW","2892.TW","5876.TW"],
        }

        # 模糊比對
        import yfinance as yf
        matched_key = None
        for key in concepts:
            if key in theme or theme in key:
                matched_key = key
                break

        if not matched_key:
            # 用 DuckDuckGo 搜尋
            try:
                from ddgs import DDGS
                with DDGS() as ddgs:
                    results = list(ddgs.text(f"台股 {theme} 概念股 相關股票", region="zh-tw", max_results=5))
                lines = [f"🏭 {theme} 概念股\n（搜尋結果）"]
                for r in results:
                    lines.append(f"• {r.get('title','')}\n  {r.get('body','')[:100]}")
                return "\n\n".join(lines)
            except Exception:
                return f"找不到「{theme}」概念股，支援：{'、'.join(concepts.keys())}"

        syms = concepts[matched_key]
        lines = [f"🏭 {matched_key} 概念股（台股）\n"]
        for sym in syms:
            try:
                hist = yf.Ticker(sym).history(period="2d")
                info = yf.Ticker(sym).info
                name = info.get("shortName") or sym.replace(".TW","")
                if len(hist) >= 2:
                    chg = (hist["Close"].iloc[-1] / hist["Close"].iloc[-2] - 1) * 100
                    price = hist["Close"].iloc[-1]
                    arrow = "▲" if chg >= 0 else "▼"
                    lines.append(f"{sym.replace('.TW','')} {name}：{price:.1f}  {arrow}{chg:+.1f}%")
            except Exception:
                lines.append(sym.replace(".TW",""))
        return "\n".join(lines)
    except Exception as e:
        return f"概念股查詢失敗：{e}"


def fetch_correlation(symbols: list, period: str = "1y") -> str:
    try:
        import yfinance as yf
        import pandas as pd
        data = {}
        for sym in symbols[:5]:
            hist = yf.Ticker(sym).history(period=period)
            if not hist.empty:
                data[sym] = hist["Close"].pct_change().dropna()

        if len(data) < 2:
            return "至少需要 2 個有效的股票代號"

        df = pd.DataFrame(data).dropna()
        corr = df.corr()

        lines = [f"📊 相關性矩陣（{period}）\n"]
        syms = list(corr.columns)
        header = "      " + "  ".join(f"{s:>8}" for s in syms)
        lines.append(header)
        for s1 in syms:
            row = f"{s1:>6}"
            for s2 in syms:
                val = corr.loc[s1, s2]
                row += f"  {val:>8.3f}"
            lines.append(row)

        lines.append("\n💡 相關係數：1.0=完全正相關，0=無關，-1.0=完全負相關")
        lines.append("分散風險建議選相關係數 < 0.5 的資產")

        # 找最低相關配對
        pairs = []
        for i, s1 in enumerate(syms):
            for s2 in syms[i+1:]:
                pairs.append((corr.loc[s1, s2], s1, s2))
        pairs.sort()
        if pairs:
            v, s1, s2 = pairs[0]
            lines.append(f"最低相關：{s1} & {s2}（{v:.3f}）← 最佳分散組合")
        return "\n".join(lines)
    except Exception as e:
        return f"相關性分析失敗：{e}"


def fetch_critique_writer(subject: str, type_: str = "觀點", lang: str = "zh-tw") -> str:
    try:
        from ddgs import DDGS
        region = "tw-tzh" if lang == "zh-tw" else "us-en"
        queries = [
            f"{subject} 評論 分析 優點",
            f"{subject} 批評 缺點 問題 盲點",
            f"{subject} 背景 脈絡 假設",
        ]
        sections = {}
        with DDGS() as ddgs:
            for q in queries:
                hits = list(ddgs.text(q, region=region, max_results=3))
                sections[q] = hits
        lines = [f"🔬 批判性評析：{subject}（{type_}）\n"]
        lines.append("── 優點與貢獻 ──")
        for h in sections[queries[0]][:2]:
            lines.append(f"✅ {h.get('title','')}：{h.get('body','')[:180]}")
        lines.append("\n── 盲點與缺失 ──")
        for h in sections[queries[1]][:2]:
            lines.append(f"⚠️ {h.get('title','')}：{h.get('body','')[:180]}")
        lines.append("\n── 隱藏假設與脈絡 ──")
        for h in sections[queries[2]][:2]:
            lines.append(f"🔍 {h.get('title','')}：{h.get('body','')[:180]}")
        lines += [
            "\n── 改進建議 ──",
            "（Claude 將根據以上資料，在回應中提出具體改進方向）",
            "\n💡 好的批判不是否定，而是幫助對象看見自己看不見的角落。",
        ]
        return "\n".join(lines)
    except Exception as e:
        return f"批判性評析失敗：{e}"


def fetch_crypto(coin: str, vs_currency: str = "usd") -> str:
    try:
        ticker_map = {
            "btc": "bitcoin", "eth": "ethereum", "sol": "solana",
            "bnb": "binancecoin", "xrp": "ripple", "doge": "dogecoin",
            "ada": "cardano", "dot": "polkadot", "matic": "matic-network",
            "avax": "avalanche-2", "link": "chainlink", "uni": "uniswap",
            "ltc": "litecoin", "bch": "bitcoin-cash", "atom": "cosmos",
            "trx": "tron", "etc": "ethereum-classic", "shib": "shiba-inu",
        }
        coin_id = ticker_map.get(coin.lower(), coin.lower())
        url = f"https://api.coingecko.com/api/v3/coins/{coin_id}?localization=false&tickers=false&community_data=false&developer_data=false"
        resp = requests.get(url, timeout=10)
        data = resp.json()
        if "error" in data:
            return f"找不到幣種「{coin}」，請確認名稱"
        md = data["market_data"]
        name = data["name"]
        sym = data["symbol"].upper()
        cur = vs_currency
        price = md["current_price"][cur]
        ch24 = md["price_change_percentage_24h"] or 0
        ch7d = md["price_change_percentage_7d"] or 0
        mcap = md["market_cap"][cur]
        vol = md["total_volume"][cur]
        hi24 = md["high_24h"][cur]
        lo24 = md["low_24h"][cur]
        ath = md["ath"][cur]
        ath_chg = md["ath_change_percentage"][cur] or 0
        cur_label = cur.upper()
        arrow = "▲" if ch24 >= 0 else "▼"
        mc_str = f"{mcap/1e12:.2f}T" if mcap >= 1e12 else f"{mcap/1e9:.2f}B"
        return (
            f"🪙 {name} ({sym})\n"
            f"💰 現價：{price:,.4f} {cur_label}  {arrow} {abs(ch24):.2f}% (24h)\n"
            f"📅 7日漲跌：{ch7d:+.2f}%\n"
            f"📊 24h 高低：{lo24:,.4f} ~ {hi24:,.4f}\n"
            f"💎 市值：{mc_str} {cur_label}\n"
            f"📦 24h 交易量：{vol/1e9:.2f}B {cur_label}\n"
            f"🏔 歷史高點：{ath:,.4f}（距高 {ath_chg:.1f}%）"
        )
    except Exception as e:
        return f"查詢加密幣「{coin}」失敗：{e}"


def fetch_crypto_depth(coin: str = "bitcoin") -> str:
    try:
        coin_map = {
            "btc": "bitcoin", "eth": "ethereum", "sol": "solana",
            "bnb": "binancecoin", "xrp": "ripple", "doge": "dogecoin",
        }
        coin_id = coin_map.get(coin.lower(), coin.lower())

        # CoinGecko 詳細資料
        resp = requests.get(
            f"https://api.coingecko.com/api/v3/coins/{coin_id}"
            "?localization=false&tickers=false&community_data=true&developer_data=false",
            timeout=10
        )
        data = resp.json()
        if "error" in data:
            return f"找不到幣種「{coin}」"

        md = data["market_data"]
        name = data["name"]
        sym = data["symbol"].upper()
        price = md["current_price"]["usd"]
        ch24 = md.get("price_change_percentage_24h") or 0
        ch7d = md.get("price_change_percentage_7d") or 0
        mcap = md["market_cap"]["usd"]
        vol = md["total_volume"]["usd"]
        dom = data.get("market_cap_percentage", {}).get(coin_id.split("-")[0], 0)

        # 資金費率（用 DuckDuckGo 補充）
        funding_note = ""
        try:
            from ddgs import DDGS
            with DDGS() as ddgs:
                r = list(ddgs.text(f"{sym} funding rate perpetual", max_results=1))
                if r:
                    funding_note = f"\n資金費率參考：{r[0].get('body','')[:80]}"
        except Exception:
            pass

        return (
            f"🔗 {name}（{sym}）鏈上深度\n\n"
            f"現價：${price:,.4f}\n"
            f"24h：{ch24:+.2f}%　7d：{ch7d:+.2f}%\n"
            f"市值：${mcap/1e9:.2f}B\n"
            f"24h 成交量：${vol/1e9:.2f}B\n"
            f"社群分數：{data.get('community_score',0):.0f}/100\n"
            f"開發者分數：{data.get('developer_score',0):.0f}/100"
            + funding_note
        )
    except Exception as e:
        return f"加密幣深度查詢失敗：{e}"


def fetch_currency_converter(amount: float, from_currency: str, to_currency: str) -> str:
    try:
        import yfinance as yf
        fc = from_currency.upper()
        tc = to_currency.upper()
        if fc == tc:
            return f"{amount:,.2f} {fc} = {amount:,.2f} {tc}"
        # 嘗試直接查匯率
        pair = f"{fc}{tc}=X"
        ticker = yf.Ticker(pair)
        hist = ticker.history(period="2d")
        if hist.empty:
            # 嘗試反向
            pair2 = f"{tc}{fc}=X"
            hist2 = yf.Ticker(pair2).history(period="2d")
            if hist2.empty:
                return f"無法取得 {fc}/{tc} 匯率資料"
            rate = 1 / hist2["Close"].iloc[-1]
        else:
            rate = hist["Close"].iloc[-1]
        result = amount * rate
        lines = [
            f"💱 外幣換算\n",
            f"{amount:,.2f} {fc}",
            f"= {result:,.4f} {tc}",
            f"\n即時匯率：1 {fc} = {rate:.4f} {tc}",
        ]
        return "\n".join(lines)
    except Exception as e:
        return f"外幣換算失敗：{e}"


def fetch_debate_simulator(motion: str, lang: str = "zh-tw") -> str:
    try:
        from ddgs import DDGS
        region = "tw-tzh" if lang == "zh-tw" else "us-en"
        pro_hits, con_hits = [], []
        with DDGS() as ddgs:
            pro_hits = list(ddgs.text(f"{motion} 支持 贊成 優點 好處", region=region, max_results=4))
            con_hits = list(ddgs.text(f"{motion} 反對 質疑 缺點 問題", region=region, max_results=4))
        lines = [
            f"⚔️ 辯論模擬：{motion}\n",
            f"══ 正方論點 ══",
        ]
        for h in pro_hits[:3]:
            lines.append(f"✅ {h.get('title','')}：{h.get('body','')[:180]}")
        lines += [f"\n══ 反方論點 ══"]
        for h in con_hits[:3]:
            lines.append(f"❌ {h.get('title','')}：{h.get('body','')[:180]}")
        pro_strength = len(pro_hits)
        con_strength = len(con_hits)
        if pro_strength > con_strength:
            verdict = "正方論據較充分"
        elif con_strength > pro_strength:
            verdict = "反方論據較充分"
        else:
            verdict = "雙方論據相當，難以定論"
        lines += [
            f"\n══ 綜合判斷 ══",
            f"議題：「{motion}」",
            f"評估：{verdict}",
            f"此議題涉及多方面考量，建議結合個人價值觀與具體情境做判斷。",
        ]
        return "\n".join(lines)
    except Exception as e:
        return f"辯論模擬失敗：{e}"


def fetch_decision_helper(question: str, options: list = None, criteria: list = None) -> str:
    try:
        from ddgs import DDGS
        lines = [f"🤔 決策輔助：{question}\n"]
        if options:
            lines.append(f"選項：{' vs '.join(options)}\n")
        if criteria:
            lines.append(f"考量標準：{', '.join(criteria)}\n")
        # 搜尋相關資訊
        with DDGS() as ddgs:
            hits = list(ddgs.text(f"{question} 建議 分析 怎麼決定", region="tw-tzh", max_results=5))
        lines.append("── 相關資訊 ──")
        for h in hits[:4]:
            lines.append(f"• {h.get('title','')}：{h.get('body','')[:200]}")
        lines.append("")
        if options:
            lines.append("── 各選項分析 ──")
            with DDGS() as ddgs:
                for opt in options[:3]:
                    opt_hits = list(ddgs.text(f"{opt} 優缺點 評價", region="tw-tzh", max_results=2))
                    lines.append(f"【{opt}】")
                    for h in opt_hits[:1]:
                        lines.append(f"  {h.get('body','')[:180]}")
                    lines.append("")
        lines.append("── 建議框架 ──")
        if criteria:
            for c in criteria:
                lines.append(f"□ {c}：請根據您的具體情況評分")
        lines.append("\n綜合以上資訊，建議依個人優先順序做最終判斷。")
        return "\n".join(lines)
    except Exception as e:
        return f"決策輔助失敗：{e}"


def fetch_deep_research(topic: str, lang: str = "zh-tw", depth: int = 5) -> str:
    try:
        from ddgs import DDGS
        depth = int(depth) if not isinstance(depth, int) else depth
        depth = min(max(depth, 3), 8)
        # 自動生成子問題
        sub_questions = [
            f"{topic} 是什麼 基本介紹",
            f"{topic} 最新發展 現況",
            f"{topic} 數據 統計 報告",
            f"{topic} 爭議 問題 缺點",
            f"{topic} 未來趨勢 預測",
            f"{topic} 專家看法 分析",
            f"{topic} 影響 重要性",
            f"{topic} 解決方案 建議",
        ][:depth]
        results = {}
        with DDGS() as ddgs:
            for q in sub_questions:
                hits = list(ddgs.text(q, region="tw-tzh" if lang == "zh-tw" else "us-en", max_results=3))
                if hits:
                    results[q] = hits
        lines = [f"📚 深度研究：{topic}\n"]
        for q, hits in results.items():
            lines.append(f"【{q}】")
            for h in hits:
                title = h.get("title", "")
                body = h.get("body", "")[:200]
                lines.append(f"  • {title}：{body}")
            lines.append("")
        lines.append(f"共蒐集 {len(results)} 個面向，{sum(len(v) for v in results.values())} 筆資料")
        return "\n".join(lines)
    except Exception as e:
        return f"深度研究失敗：{e}"


def fetch_defi_calculator(principal_usd: float, apy: float, days: int,
                           compound: bool = True, protocol: str = "") -> str:
    try:
        import yfinance as yf
        # 查詢USD/TWD匯率
        try:
            usdtwd = yf.Ticker("USDTWD=X").history(period="2d")["Close"].iloc[-1]
        except Exception:
            usdtwd = 32.0
        if compound:
            daily_rate = apy / 100 / 365
            final_usd = principal_usd * (1 + daily_rate) ** days
        else:
            final_usd = principal_usd * (1 + apy / 100 * days / 365)
        profit_usd = final_usd - principal_usd
        lines = [
            f"🔗 DeFi收益試算\n",
        ]
        if protocol:
            lines.append(f"協議：{protocol}")
        lines += [
            f"本金：${principal_usd:,.2f} USD（約 {principal_usd*usdtwd:,.0f} TWD）",
            f"APY：{apy}%　質押天數：{days} 天",
            f"複利：{'是' if compound else '否'}",
            f"",
            f"── 試算結果 ──",
            f"到期本利和：${final_usd:,.4f} USD",
            f"獲利：${profit_usd:,.4f} USD（約 {profit_usd*usdtwd:,.0f} TWD）",
            f"實際年化（{days}天）：{(final_usd/principal_usd-1)*365/days*100:.2f}%",
            f"",
            f"⚠️ DeFi有智能合約風險、無常損失風險，本試算僅供參考",
        ]
        return "\n".join(lines)
    except Exception as e:
        return f"DeFi試算失敗：{e}"


def fetch_devil_advocate(position: str, lang: str = "zh-tw") -> str:
    try:
        from ddgs import DDGS
        region = "tw-tzh" if lang == "zh-tw" else "us-en"
        queries = [
            f"{position} 反對 批評 缺點 問題",
            f"{position} 失敗案例 風險 危險",
            f"{position} 質疑 挑戰 反駁",
        ]
        all_hits = []
        with DDGS() as ddgs:
            for q in queries:
                hits = list(ddgs.text(q, region=region, max_results=3))
                all_hits.extend(hits)
        lines = [
            f"😈 魔鬼代言人：挑戰「{position}」\n",
            f"以下從相反角度提出最強反駁，幫助您找出盲點：\n",
            f"── 反駁論點 ──",
        ]
        seen = set()
        for h in all_hits[:6]:
            title = h.get("title","")
            if title in seen:
                continue
            seen.add(title)
            body = h.get("body","")[:200]
            lines.append(f"⚡ {title}：{body}\n")
        lines.append("── 結語 ──")
        lines.append("以上為刻意的反面論點，目的是強化您的思考。若能反駁以上論點，您的立場將更為穩固。")
        return "\n".join(lines)
    except Exception as e:
        return f"魔鬼代言人失敗：{e}"


def fetch_dividend_calendar(symbol: str) -> str:
    """除權息資訊"""
    try:
        import yfinance as yf
        ticker = yf.Ticker(symbol)
        info = ticker.info
        name = info.get("longName") or info.get("shortName") or symbol
        currency = info.get("currency", "")

        div_yield    = info.get("dividendYield", 0) or 0
        div_rate     = info.get("dividendRate", 0) or 0
        ex_date      = info.get("exDividendDate")
        last_div     = info.get("lastDividendValue", 0) or 0
        payout_ratio = info.get("payoutRatio", 0) or 0
        price        = info.get("currentPrice") or info.get("regularMarketPrice", 0)

        lines = [f"💰 {name} ({symbol}) 除權息資訊\n"]
        if div_rate:
            lines.append(f"年配息：{div_rate:.4f} {currency}")
        if div_yield:
            lines.append(f"殖利率：{div_yield*100:.2f}%")
        if last_div:
            lines.append(f"上次配息：{last_div:.4f} {currency}")
        if ex_date:
            import datetime
            ex_dt = dt.datetime.fromtimestamp(ex_date).strftime("%Y-%m-%d")
            lines.append(f"除息日：{ex_dt}")
        if payout_ratio:
            lines.append(f"配息率：{payout_ratio*100:.1f}%")

        # 歷史配息紀錄
        try:
            divs = ticker.dividends
            if divs is not None and not divs.empty:
                lines.append("\n近5次配息：")
                for dt, val in divs.tail(5).iloc[::-1].items():
                    lines.append(f"  {str(dt)[:10]}：{val:.4f} {currency}")
        except Exception:
            pass

        return "\n".join(lines) if len(lines) > 1 else f"{symbol} 無配息資料"
    except Exception as e:
        return f"除權息查詢失敗：{e}"


def fetch_drag_drop(from_x: int = None, from_y: int = None, to_x: int = None, to_y: int = None,
                    from_text: str = "", to_text: str = "", monitor: int = 1, duration: float = 0.5) -> str:
    try:
        import mss, ctypes as _c, ctypes.wintypes as _w, time as _t, re

        def get_abs(x, y):
            with mss.mss() as sct:
                mon = sct.monitors[monitor] if monitor < len(sct.monitors) else sct.monitors[1]
                return mon["left"] + x, mon["top"] + y

        # 起點
        if from_text:
            r = fetch_vision_locate(from_text, monitor, "locate_only")
            m = re.search(r'\((-?\d+), (-?\d+)\)', r)
            if m: fx, fy = int(m.group(1)), int(m.group(2))
            else: return f"拖曳起點找不到「{from_text}」：{r}"
        else:
            fx, fy = get_abs(from_x or 0, from_y or 0)

        # 終點
        if to_text:
            r2 = fetch_vision_locate(to_text, monitor, "locate_only")
            m2 = re.search(r'\((-?\d+), (-?\d+)\)', r2)
            if m2: tx, ty = int(m2.group(1)), int(m2.group(2))
            else: return f"拖曳終點找不到「{to_text}」：{r2}"
        else:
            tx, ty = get_abs(to_x or 0, to_y or 0)

        # SendInput 拖曳（支援螢幕2負座標）
        vl, vt, vw, vh, is_phys, sx, sy = _get_virtual_desktop()
        u32 = _c.windll.user32
        class MI(_c.Structure):
            _fields_ = [('dx',_w.LONG),('dy',_w.LONG),('mouseData',_w.DWORD),
                        ('dwFlags',_w.DWORD),('time',_w.DWORD),('dwExtraInfo',_c.POINTER(_c.c_ulong))]
        class U(_c.Union): _fields_ = [('mi', MI)]
        class INP(_c.Structure):
            _anonymous_ = ('u',); _fields_ = [('type',_w.DWORD),('u',U)]
        def _send(flags, dx=0, dy=0):
            i = INP(0, U(mi=MI(dx, dy, 0, flags, 0, None)))
            u32.SendInput(1, _c.byref(i), _c.sizeof(i))
        def _norm(lx, ly):
            px = round(lx * sx); py = round(ly * sy)
            return int((px-vl)*65535//vw), int((py-vt)*65535//vh)

        fnx, fny = _norm(fx, fy)
        tnx, tny = _norm(tx, ty)
        steps = max(10, int(duration * 60))
        _send(0x0001|0x8000|0x4000, fnx, fny); _t.sleep(0.1)
        _send(0x0002)  # LEFTDOWN
        for i in range(1, steps+1):
            ix = fnx + (tnx - fnx) * i // steps
            iy = fny + (tny - fny) * i // steps
            _send(0x0001|0x8000|0x4000, ix, iy); _t.sleep(duration / steps)
        _send(0x0004)  # LEFTUP
        return f"✅ 拖曳完成：({fx},{fy}) → ({tx},{ty})，耗時 {duration}s"
    except Exception as e:
        return f"拖曳失敗：{e}"


def fetch_drip_calculator(symbol: str, shares: float, years: int = 10, monthly_invest: float = 0) -> str:
    try:
        import yfinance as yf
        ticker = yf.Ticker(symbol)
        info = ticker.info
        name = info.get("shortName") or symbol
        price = info.get("currentPrice") or info.get("regularMarketPrice") or 0
        div_yield = (info.get("dividendYield") or 0)
        div_rate = info.get("dividendRate") or (price * div_yield)
        payout_freq = 4  # 假設季配

        if not price:
            return f"找不到 {symbol} 的股價資料"

        total_shares = shares
        total_invested = shares * price
        annual_divs = []

        for year in range(1, years + 1):
            # 年度股息 → 再買股
            annual_div = total_shares * div_rate
            new_shares = annual_div / price if price > 0 else 0
            total_shares += new_shares

            # 每月定期追加
            if monthly_invest > 0:
                monthly_shares = (monthly_invest * 12) / price
                total_shares += monthly_shares
                total_invested += monthly_invest * 12

            annual_divs.append(annual_div)

        final_value = total_shares * price
        total_div = sum(annual_divs)
        total_return = (final_value / total_invested - 1) * 100 if total_invested > 0 else 0

        lines = [
            f"💰 {name}（{symbol}）DRIP 股息再投資試算\n",
            f"初始：{shares:.0f} 股 × {price:.2f} = {shares*price:,.0f} 元",
        ]
        if monthly_invest > 0:
            lines.append(f"每月追加：{monthly_invest:,.0f} 元")
        lines += [
            f"殖利率：{div_yield*100:.2f}%　每股年配：{div_rate:.4f}",
            f"\n{years} 年後：",
            f"持股數：{total_shares:,.1f} 股",
            f"總市值：{final_value:,.0f} 元",
            f"總投入：{total_invested:,.0f} 元",
            f"累積股息：{total_div:,.0f} 元",
            f"總報酬：{total_return:+.1f}%",
            f"年配息（第{years}年）：{annual_divs[-1]:,.0f} 元",
        ]
        return "\n".join(lines)
    except Exception as e:
        return f"DRIP 試算失敗：{e}"


def fetch_earnings(symbol: str) -> str:
    try:
        import yfinance as yf
        ticker = yf.Ticker(symbol)
        info = ticker.info
        name = info.get("longName") or info.get("shortName") or symbol

        lines = [f"📊 {name} ({symbol}) 財報趨勢\n"]

        # 季度財報（用 quarterly_income_stmt，quarterly_earnings 已棄用）
        try:
            qi = ticker.quarterly_income_stmt
            if qi is not None and not qi.empty and "Basic EPS" in qi.index:
                lines.append("── 近幾季 EPS ──")
                eps_row = qi.loc["Basic EPS"]
                for col in eps_row.index[:6]:
                    val = eps_row[col]
                    if val is not None and str(val) != "nan":
                        quarter = str(col)[:10]
                        lines.append(f"  {quarter}：EPS {float(val):.2f}")
        except Exception:
            pass

        # 年度財報趨勢
        try:
            fin = ticker.financials
            if fin is not None and not fin.empty:
                lines.append("\n── 年度財務 ──")
                rev_row = fin.loc["Total Revenue"] if "Total Revenue" in fin.index else None
                ni_row = fin.loc["Net Income"] if "Net Income" in fin.index else None
                cols = fin.columns[:4]
                if rev_row is not None:
                    vals = [f"{rev_row[c]/1e9:.1f}B" for c in cols if c in rev_row.index]
                    lines.append(f"  營收：{' → '.join(vals)}")
                if ni_row is not None:
                    vals = [f"{ni_row[c]/1e9:.1f}B" for c in cols if c in ni_row.index]
                    lines.append(f"  淨利：{' → '.join(vals)}")
        except Exception:
            pass

        # 下次財報日
        try:
            cal = ticker.calendar
            if cal is not None:
                ed = cal.get("Earnings Date") or cal.get("earningsDate")
                if ed is not None:
                    if hasattr(ed, '__iter__'):
                        ed = list(ed)[0]
                    lines.append(f"\n下次財報日：{str(ed)[:10]}")
        except Exception:
            pass

        return "\n".join(lines) if len(lines) > 1 else f"找不到 {symbol} 的財報資料"
    except Exception as e:
        return f"財報查詢失敗：{e}"


def fetch_earnings_calendar(days: int = 7) -> str:
    try:
        import feedparser, datetime
        lines = [f"📊 未來 {days} 天財報行事曆\n"]
        url = "https://news.google.com/rss/search?q=earnings+report+quarterly+results&hl=en-US&gl=US&ceid=US:en"
        feed = feedparser.parse(url)
        count = 0
        for entry in feed.entries:
            title = entry.get("title", "").split(" - ")[0]
            pub = entry.get("published", "")[:16]
            lines.append(f"• {title}（{pub}）")
            count += 1
            if count >= 10:
                break

        # 補充重點大型股財報
        import yfinance as yf
        watchlist = ["AAPL", "MSFT", "GOOGL", "AMZN", "NVDA", "META", "TSLA"]
        lines.append("\n重點股財報日：")
        for sym in watchlist:
            try:
                cal = yf.Ticker(sym).calendar
                if cal and "Earnings Date" in cal:
                    ed = cal["Earnings Date"]
                    if hasattr(ed, '__iter__'):
                        ed = list(ed)[0]
                    lines.append(f"  {sym}：{str(ed)[:10]}")
            except Exception:
                pass
        return "\n".join(lines)
    except Exception as e:
        return f"財報日曆查詢失敗：{e}"


def fetch_economic_calendar(count: int = 10) -> str:
    try:
        import feedparser
        results = []
        # Investing.com RSS 與 Google 新聞
        urls = [
            ("https://news.google.com/rss/search?q=CPI+非農+Fed利率+GDP+經濟數據&hl=zh-Hant&gl=TW&ceid=TW:zh-Hant", "財經日曆"),
            ("https://news.google.com/rss/search?q=economic+calendar+CPI+nonfarm+Fed+GDP&hl=en-US&gl=US&ceid=US:en", "Economic Calendar"),
        ]
        lines = ["📅 重要經濟數據行事曆\n"]
        seen = set()
        for url, label in urls:
            feed = feedparser.parse(url)
            for entry in feed.entries:
                title = entry.get("title", "").split(" - ")[0]
                pub = entry.get("published", "")[:16]
                key = title[:30]
                if key not in seen:
                    seen.add(key)
                    lines.append(f"• {title}（{pub}）")
                if len(lines) > count + 1:
                    break
            if len(lines) > count + 1:
                break
        lines.append("\n💡 重點關注：CPI（通膨）、非農就業（NFP）、Fed利率決議、GDP、PPI")
        return "\n".join(lines)
    except Exception as e:
        return f"經濟日曆查詢失敗：{e}"


def fetch_etf(symbol: str) -> str:
    try:
        import yfinance as yf
        ticker = yf.Ticker(symbol)
        info = ticker.info
        name = info.get("longName") or info.get("shortName") or symbol
        currency = info.get("currency", "")

        hist = ticker.history(period="1y")
        if hist.empty:
            return f"找不到 ETF「{symbol}」"

        current = hist["Close"].iloc[-1]
        ret_1m = ((hist["Close"].iloc[-1] / hist["Close"].iloc[-22]) - 1) * 100 if len(hist) > 22 else None
        ret_3m = ((hist["Close"].iloc[-1] / hist["Close"].iloc[-66]) - 1) * 100 if len(hist) > 66 else None
        ret_1y = ((hist["Close"].iloc[-1] / hist["Close"].iloc[0]) - 1) * 100 if len(hist) > 5 else None

        expense = info.get("annualReportExpenseRatio")
        div_yield = info.get("dividendYield") or info.get("yield")
        total_assets = info.get("totalAssets")
        category = info.get("category", "")

        # 前10大持股
        try:
            holdings = ticker.funds_data.top_holdings
            top_str = ""
            if holdings is not None and not holdings.empty:
                top_items = []
                for _, row in holdings.head(5).iterrows():
                    n = row.get("Name") or row.get("name") or ""
                    w = row.get("Holding Percent") or row.get("holdingPercent") or row.get("weight") or 0
                    top_items.append(f"  {n}（{w*100:.1f}%）")
                if top_items:
                    top_str = "\n前5大持股：\n" + "\n".join(top_items)
        except Exception:
            top_str = ""

        lines = [f"📦 {name} ({symbol})\n現價：{current:.2f} {currency}\n"]
        if category: lines.append(f"類型：{category}")
        if total_assets: lines.append(f"規模：{total_assets/1e9:.1f}B {currency}")
        if expense: lines.append(f"費用率：{expense*100:.2f}%")
        if div_yield: lines.append(f"配息殖利率：{div_yield*100:.2f}%")
        lines.append("\n── 績效 ──")
        if ret_1m: lines.append(f"近1月：{ret_1m:+.2f}%")
        if ret_3m: lines.append(f"近3月：{ret_3m:+.2f}%")
        if ret_1y: lines.append(f"近1年：{ret_1y:+.2f}%")
        if top_str: lines.append(top_str)

        return "\n".join(lines)
    except Exception as e:
        return f"ETF 查詢失敗：{e}"


def fetch_fact_check(claim: str, lang: str = "zh-tw") -> str:
    try:
        from ddgs import DDGS
        region = "tw-tzh" if lang == "zh-tw" else "us-en"
        queries = [
            f"{claim} 真假 查核",
            f"{claim} 事實查核",
            f"{claim} 錯誤 謠言",
            f"{claim} 正確 證實",
        ]
        all_hits = []
        with DDGS() as ddgs:
            for q in queries:
                hits = list(ddgs.text(q, region=region, max_results=3))
                all_hits.extend(hits)
        # 簡單關鍵字判斷情緒
        text_blob = " ".join(h.get("title","") + " " + h.get("body","") for h in all_hits).lower()
        false_kw = ["假", "謠言", "錯誤", "誤導", "不實", "false", "fake", "wrong", "misleading"]
        true_kw = ["真", "確認", "屬實", "正確", "true", "confirmed", "correct", "verified"]
        false_score = sum(text_blob.count(k) for k in false_kw)
        true_score = sum(text_blob.count(k) for k in true_kw)
        if false_score > true_score * 2:
            verdict = "❌ 可能為假/誤導"
        elif true_score > false_score * 2:
            verdict = "✅ 可能屬實"
        elif false_score > 0 or true_score > 0:
            verdict = "⚠️ 有爭議，需進一步確認"
        else:
            verdict = "❓ 資料不足，無法判斷"
        lines = [f"🔍 事實查核\n", f"說法：「{claim}」\n", f"查核結果：{verdict}\n", f"── 相關資料 ──"]
        for h in all_hits[:6]:
            title = h.get("title", "")
            body = h.get("body", "")[:150]
            url = h.get("href", "")
            lines.append(f"• {title}\n  {body}")
        return "\n".join(lines)
    except Exception as e:
        return f"事實查核失敗：{e}"


def fetch_finance_news(source: str = "all", count: int = 5) -> str:
    try:
        import feedparser
        count = min(max(count, 1), 10)
        feeds = {
            "yahoo_tw": ("Yahoo奇摩財經", "https://tw.stock.yahoo.com/rss"),
            "cnyes":    ("鉅亨網",        "https://feeds.feedburner.com/cnyes"),
            "udn":      ("聯合財經網",    "https://udn.com/rssfeed/news/2/6644?ch=news"),
            "moneydj":  ("MoneyDJ",      "https://www.moneydj.com/KMDJ/RssNew/RssNewList.aspx?index=1&param="),
            "ctee":     ("工商時報",      "https://www.ctee.com.tw/feeds/latest"),
            "yahoo_us": ("Yahoo Finance", "https://finance.yahoo.com/news/rssindex"),
        }
        sources = list(feeds.keys()) if source == "all" else [source]
        results = []
        for src in sources:
            if src not in feeds:
                continue
            label, url = feeds[src]
            try:
                feed = feedparser.parse(url)
                items = feed.entries[:count]
                if not items:
                    results.append(f"📰 {label}：暫無資料")
                    continue
                lines = [f"📰 {label}"]
                for i, entry in enumerate(items, 1):
                    title = entry.get("title", "無標題")
                    lines.append(f"{i}. {title}")
                results.append("\n".join(lines))
            except Exception:
                results.append(f"📰 {label}：抓取失敗")
        return "\n\n".join(results) if results else "無法取得新聞"
    except Exception as e:
        return f"財經新聞失敗：{e}"


def fetch_financial_health(monthly_income: float, monthly_expense: float,
                            total_assets: float, total_debt: float,
                            emergency_fund_months: float = 0,
                            has_insurance: bool = False,
                            investment_ratio: float = 0) -> str:
    try:
        score = 100
        issues = []
        goods = []
        # 儲蓄率
        save_rate = (monthly_income - monthly_expense) / monthly_income * 100 if monthly_income > 0 else 0
        if save_rate < 0:
            score -= 30; issues.append("每月支出超過收入（負儲蓄）")
        elif save_rate < 10:
            score -= 15; issues.append(f"儲蓄率偏低（{save_rate:.1f}%，建議≥20%）")
        elif save_rate >= 20:
            goods.append(f"儲蓄率良好（{save_rate:.1f}%）")
        # 負債比
        debt_ratio = total_debt / total_assets * 100 if total_assets > 0 else 100
        if debt_ratio > 70:
            score -= 25; issues.append(f"負債比過高（{debt_ratio:.1f}%，建議＜50%）")
        elif debt_ratio > 50:
            score -= 10; issues.append(f"負債比偏高（{debt_ratio:.1f}%）")
        else:
            goods.append(f"負債比健康（{debt_ratio:.1f}%）")
        # 緊急備用金
        if emergency_fund_months < 3:
            score -= 20; issues.append(f"緊急備用金不足（{emergency_fund_months}個月，建議≥6個月）")
        elif emergency_fund_months >= 6:
            goods.append(f"緊急備用金充足（{emergency_fund_months}個月）")
        # 保險
        if not has_insurance:
            score -= 10; issues.append("缺乏壽險/重疾險保障")
        else:
            goods.append("有保險保障")
        # 投資比例
        if investment_ratio >= 20:
            goods.append(f"積極投資（{investment_ratio}%收入）")
        elif investment_ratio > 0:
            issues.append(f"投資比例偏低（{investment_ratio}%，建議≥20%）")
            score -= 5
        # 評分
        score = max(0, min(100, score))
        if score >= 80:
            grade, emoji = "優良", "🟢"
        elif score >= 60:
            grade, emoji = "尚可", "🟡"
        elif score >= 40:
            grade, emoji = "需改善", "🟠"
        else:
            grade, emoji = "高風險", "🔴"
        lines = [
            f"💊 財務健康診斷\n",
            f"月收入：{monthly_income:,.0f}　月支出：{monthly_expense:,.0f}",
            f"總資產：{total_assets/10000:,.0f}萬　總負債：{total_debt/10000:,.0f}萬",
            f"",
            f"── 健康評分 ──",
            f"{emoji} {score} 分 / 100（{grade}）",
            f"儲蓄率：{save_rate:.1f}%　負債比：{debt_ratio:.1f}%",
        ]
        if goods:
            lines.append("\n✅ 優點：" + "、".join(goods))
        if issues:
            lines.append("\n⚠️ 待改善：")
            for i in issues:
                lines.append(f"  • {i}")
        return "\n".join(lines)
    except Exception as e:
        return f"財務健康診斷失敗：{e}"


def fetch_forex(pair: str) -> str:
    try:
        import yfinance as yf
        ticker = yf.Ticker(pair)
        hist = ticker.history(period="5d")
        if hist.empty:
            return f"找不到匯率「{pair}」，請確認格式（如 USDTWD=X）"
        info = ticker.info
        current = hist["Close"].iloc[-1]
        prev = hist["Close"].iloc[-2] if len(hist) > 1 else current
        change = current - prev
        change_pct = (change / prev * 100) if prev else 0
        arrow = "▲" if change >= 0 else "▼"
        name = info.get("shortName", pair)
        return (
            f"💱 {name}\n"
            f"匯率：{current:.4f}  {arrow} {abs(change):.4f} ({change_pct:+.2f}%)\n"
            f"近5日高：{hist['High'].max():.4f}\n"
            f"近5日低：{hist['Low'].min():.4f}"
        )
    except Exception as e:
        return f"查詢匯率「{pair}」失敗：{e}"


def fetch_forex_chart(pair: str, period: str = "3mo") -> str:
    try:
        import yfinance as yf
        ticker = yf.Ticker(pair)
        hist = ticker.history(period="6mo")
        info = ticker.info

        if hist.empty:
            return f"找不到匯率「{pair}」，請確認格式（如 USDTWD=X）"

        name = info.get("shortName") or pair
        close = hist["Close"]
        current = close.iloc[-1]
        prev = close.iloc[-2] if len(close) > 1 else current
        chg = (current / prev - 1) * 100

        ma5  = close.tail(5).mean()
        ma20 = close.tail(20).mean()
        ma60 = close.tail(60).mean() if len(close) >= 60 else close.mean()
        rsi  = calc_rsi(close) if len(close) >= 15 else None

        if ma5 > ma20 > ma60:   trend = "強勢升值 📈"
        elif ma5 < ma20 < ma60: trend = "強勢貶值 📉"
        elif ma5 > ma20:        trend = "短線偏強 🔼"
        else:                   trend = "短線偏弱 🔽"

        from ta.volatility import BollingerBands
        bb = BollingerBands(close)
        bb_upper = bb.bollinger_hband().iloc[-1]
        bb_lower = bb.bollinger_lband().iloc[-1]
        bb_pos = "近上軌（偏強）" if current >= bb_upper * 0.99 else "近下軌（偏弱）" if current <= bb_lower * 1.01 else "通道中間"

        result = (
            f"💱 {name} 技術分析\n"
            f"現值：{current:.4f}  {'▲' if chg >= 0 else '▼'} {chg:+.2f}%\n\n"
            f"MA5：{ma5:.4f}　MA20：{ma20:.4f}　MA60：{ma60:.4f}\n"
            f"趨勢：{trend}\n"
            f"布林：{bb_pos}\n"
        )
        if rsi:
            rsi_note = "超買" if rsi >= 70 else "超賣" if rsi <= 30 else "中性"
            result += f"RSI(14)：{rsi}（{rsi_note}）\n"
        result += (
            f"\n52週高：{hist['High'].max():.4f}\n"
            f"52週低：{hist['Low'].min():.4f}"
        )
        return result
    except Exception as e:
        return f"外匯技術分析失敗：{e}"


def fetch_forex_deposit(amount_twd: float, currency: str, annual_rate: float, months: int,
                         buy_rate: float = None, sell_rate: float = None) -> str:
    try:
        import yfinance as yf
        cur = currency.upper()
        # 查詢即時匯率
        if not buy_rate:
            pair = f"TWD{cur}=X"
            hist = yf.Ticker(pair).history(period="2d")
            if hist.empty:
                pair2 = f"{cur}TWD=X"
                hist2 = yf.Ticker(pair2).history(period="2d")
                if not hist2.empty:
                    rate_per_twd = 1 / hist2["Close"].iloc[-1]
                else:
                    return f"無法取得 {cur}/TWD 匯率"
            else:
                rate_per_twd = hist["Close"].iloc[-1]
            buy_rate_actual = 1 / rate_per_twd  # TWD→外幣
        else:
            rate_per_twd = 1 / buy_rate
            buy_rate_actual = buy_rate
        if not sell_rate:
            sell_rate_actual = buy_rate_actual  # 保守假設相同
        else:
            sell_rate_actual = sell_rate
        # 換算外幣本金
        foreign_principal = amount_twd / buy_rate_actual
        # 計算外幣到期本利和
        foreign_final = foreign_principal * (1 + annual_rate / 100 * months / 12)
        foreign_interest = foreign_final - foreign_principal
        # 換回台幣
        twd_final = foreign_final * sell_rate_actual
        twd_profit = twd_final - amount_twd
        effective_rate = (twd_profit / amount_twd / (months / 12) * 100) if amount_twd > 0 and months > 0 else 0
        lines = [
            f"🌐 外幣定存試算（{cur}）\n",
            f"台幣本金：{amount_twd:,.0f} 元",
            f"買入匯率：1 {cur} = {buy_rate_actual:.4f} TWD",
            f"外幣本金：{foreign_principal:,.2f} {cur}",
            f"年利率：{annual_rate}%　存款期：{months} 個月",
            f"",
            f"── 到期結果 ──",
            f"外幣本利和：{foreign_final:,.4f} {cur}",
            f"外幣利息：{foreign_interest:,.4f} {cur}",
            f"賣出匯率：1 {cur} = {sell_rate_actual:.4f} TWD",
            f"換回台幣：{twd_final:,.0f} 元",
            f"台幣獲利：{twd_profit:+,.0f} 元",
            f"等效台幣年利率：{effective_rate:.2f}%",
        ]
        if abs(buy_rate_actual - sell_rate_actual) < 0.001:
            lines.append("\n⚠️ 未考慮匯差手續費及匯率變動風險")
        return "\n".join(lines)
    except Exception as e:
        return f"外幣定存試算失敗：{e}"


def fetch_fund(symbol: str) -> str:
    try:
        import yfinance as yf
        ticker = yf.Ticker(symbol)
        info = ticker.info
        hist = ticker.history(period="1y")
        if hist.empty:
            return f"找不到 {symbol} 的基金資料"
        price = hist["Close"].iloc[-1]
        price_1m = hist["Close"].iloc[-22] if len(hist) > 22 else hist["Close"].iloc[0]
        price_3m = hist["Close"].iloc[-66] if len(hist) > 66 else hist["Close"].iloc[0]
        price_1y = hist["Close"].iloc[0]
        ret_1m = (price / price_1m - 1) * 100
        ret_3m = (price / price_3m - 1) * 100
        ret_1y = (price / price_1y - 1) * 100
        name = info.get("longName") or info.get("shortName") or symbol
        expense = info.get("annualReportExpenseRatio")
        category = info.get("category") or info.get("fundFamily") or "—"
        nav = info.get("navPrice") or price
        lines = [
            f"📦 基金查詢：{name}\n",
            f"代號：{symbol}　類別：{category}",
            f"淨值/價格：{nav:.2f}",
        ]
        if expense:
            lines.append(f"費用率：{expense*100:.2f}%")
        lines += [
            f"",
            f"── 績效 ──",
            f"近1月：{ret_1m:+.2f}%",
            f"近3月：{ret_3m:+.2f}%",
            f"近1年：{ret_1y:+.2f}%",
        ]
        return "\n".join(lines)
    except Exception as e:
        return f"基金查詢失敗：{e}"


def fetch_fundamentals(symbol: str) -> str:
    try:
        import yfinance as yf
        ticker = yf.Ticker(symbol)
        info = ticker.info
        name = info.get("longName") or info.get("shortName") or symbol
        currency = info.get("currency", "")

        # 基本面指標
        roe = info.get("returnOnEquity")
        roa = info.get("returnOnAssets")
        profit_margin = info.get("profitMargins")
        gross_margin = info.get("grossMargins")
        operating_margin = info.get("operatingMargins")
        debt_equity = info.get("debtToEquity")
        current_ratio = info.get("currentRatio")
        eps = info.get("trailingEps")
        eps_fwd = info.get("forwardEps")
        pe = info.get("trailingPE")
        pe_fwd = info.get("forwardPE")
        pb = info.get("priceToBook")
        rev_growth = info.get("revenueGrowth")
        earn_growth = info.get("earningsGrowth")
        dividend_yield = info.get("dividendYield")

        # 分析師評級
        target_mean = info.get("targetMeanPrice")
        target_high = info.get("targetHighPrice")
        target_low = info.get("targetLowPrice")
        recommend = info.get("recommendationKey", "")
        recommend_map = {"strong_buy": "強力買進 💚", "buy": "買進 🟢", "hold": "持有 🟡",
                         "sell": "賣出 🔴", "strong_sell": "強力賣出 ❌"}
        recommend_str = recommend_map.get(recommend, recommend)
        num_analysts = info.get("numberOfAnalystOpinions", 0)

        lines = [f"📋 {name} ({symbol}) 深度基本面\n"]

        lines.append("── 獲利能力 ──")
        if roe: lines.append(f"ROE（股東權益報酬）：{roe*100:.1f}%")
        if roa: lines.append(f"ROA（資產報酬）：{roa*100:.1f}%")
        if gross_margin: lines.append(f"毛利率：{gross_margin*100:.1f}%")
        if operating_margin: lines.append(f"營業利益率：{operating_margin*100:.1f}%")
        if profit_margin: lines.append(f"淨利率：{profit_margin*100:.1f}%")

        lines.append("\n── 估值 ──")
        if pe: lines.append(f"本益比（P/E）：{pe:.1f}")
        if pe_fwd: lines.append(f"預估本益比：{pe_fwd:.1f}")
        if pb: lines.append(f"股價淨值比（P/B）：{pb:.2f}")
        if eps: lines.append(f"EPS（過去12月）：{eps:.2f} {currency}")
        if eps_fwd: lines.append(f"EPS 預估：{eps_fwd:.2f} {currency}")

        lines.append("\n── 成長性 ──")
        if rev_growth: lines.append(f"營收年增率：{rev_growth*100:+.1f}%")
        if earn_growth: lines.append(f"獲利年增率：{earn_growth*100:+.1f}%")

        lines.append("\n── 財務健康 ──")
        if debt_equity: lines.append(f"負債股權比：{debt_equity:.1f}%")
        if current_ratio: lines.append(f"流動比率：{current_ratio:.2f}")
        if dividend_yield: lines.append(f"殖利率：{dividend_yield*100:.2f}%")

        if target_mean and num_analysts:
            lines.append(f"\n── 分析師（{num_analysts} 位）──")
            lines.append(f"評級：{recommend_str}")
            lines.append(f"目標價：{target_low:.2f} ~ {target_high:.2f}（均值 {target_mean:.2f} {currency}）")

        return "\n".join(lines)
    except Exception as e:
        return f"基本面查詢失敗：{e}"


def fetch_futures(items: list = None) -> str:
    """主要期貨報價"""
    try:
        import yfinance as yf
        futures_map = {
            "sp500":  ("S&P500期貨",    "ES=F"),
            "nasdaq": ("那斯達克期貨",  "NQ=F"),
            "dow":    ("道瓊期貨",      "YM=F"),
            "gold":   ("黃金期貨",      "GC=F"),
            "oil":    ("WTI原油期貨",   "CL=F"),
            "taiex":  ("台指期",        "TAIEX.TW"),
        }
        if not items or "all" in items:
            items = list(futures_map.keys())

        lines = ["📊 期貨報價\n"]
        for key in items:
            if key not in futures_map:
                continue
            name, sym = futures_map[key]
            try:
                hist = yf.Ticker(sym).history(period="2d")
                if hist.empty:
                    continue
                price = hist["Close"].iloc[-1]
                if len(hist) >= 2:
                    chg = (price / hist["Close"].iloc[-2] - 1) * 100
                    arrow = "▲" if chg >= 0 else "▼"
                    lines.append(f"{name}：{price:,.2f}  {arrow} {chg:+.2f}%")
                else:
                    lines.append(f"{name}：{price:,.2f}")
            except Exception:
                pass
        return "\n".join(lines)
    except Exception as e:
        return f"期貨查詢失敗：{e}"


def fetch_global_market() -> str:
    try:
        import yfinance as yf
        markets = {
            "🇺🇸 S&P500": "^GSPC", "🇺🇸 那斯達克": "^IXIC", "🇺🇸 道瓊": "^DJI",
            "🇹🇼 台股": "^TWII", "🇭🇰 恆生": "^HSI", "🇯🇵 日經": "^N225",
            "🇰🇷 韓股": "^KS11", "🇬🇧 英國": "^FTSE", "🇩🇪 德國": "^GDAXI",
            "🇫🇷 法國": "^FCHI", "🌏 上證": "000001.SS",
        }
        lines = ["🌍 全球市場概覽\n"]
        for name, sym in markets.items():
            try:
                hist = yf.Ticker(sym).history(period="2d")
                if len(hist) >= 2:
                    chg = (hist["Close"].iloc[-1] / hist["Close"].iloc[-2] - 1) * 100
                    price = hist["Close"].iloc[-1]
                    arrow = "▲" if chg >= 0 else "▼"
                    lines.append(f"{name}：{price:,.2f}  {arrow} {chg:+.2f}%")
            except Exception:
                pass
        return "\n".join(lines)
    except Exception as e:
        return f"全球市場查詢失敗：{e}"


def fetch_gold_calculator(weight: float, unit: str = "公克", currency: str = "TWD") -> str:
    try:
        import yfinance as yf
        # 查詢黃金價格（美元/盎司）
        gold = yf.Ticker("GC=F").history(period="2d")
        if gold.empty:
            gold = yf.Ticker("GLD").history(period="2d")
            if gold.empty:
                return "無法取得黃金價格"
            gold_usd_oz = gold["Close"].iloc[-1] / 0.0965835  # GLD每股≈0.0965835盎司
        else:
            gold_usd_oz = gold["Close"].iloc[-1]
        # 換算匯率
        try:
            usdtwd = yf.Ticker("USDTWD=X").history(period="2d")["Close"].iloc[-1]
        except Exception:
            usdtwd = 32.0
        # 單位換算為公克
        unit_map = {"公克": 1, "錢": 3.75, "兩": 37.5, "盎司": 31.1035}
        gram = weight * unit_map.get(unit, 1)
        # 公克換算盎司
        oz = gram / 31.1035
        value_usd = oz * gold_usd_oz
        value_twd = value_usd * usdtwd
        gold_per_gram_twd = gold_usd_oz / 31.1035 * usdtwd
        lines = [
            f"🥇 黃金換算\n",
            f"國際金價：${gold_usd_oz:,.2f}/盎司",
            f"每公克：{gold_per_gram_twd:,.0f} TWD　${gold_usd_oz/31.1035:.2f} USD",
            f"",
            f"── 換算結果 ──",
            f"{weight} {unit} = {gram:.4f} 公克 = {oz:.6f} 盎司",
        ]
        if currency == "TWD":
            lines.append(f"價值：{value_twd:,.0f} 新台幣（匯率 {usdtwd:.2f}）")
        else:
            lines.append(f"價值：${value_usd:,.2f} USD")
        return "\n".join(lines)
    except Exception as e:
        return f"黃金換算失敗：{e}"


def fetch_google_trends(keywords: list, timeframe: str = "today 3-m", geo: str = "TW") -> str:
    try:
        from pytrends.request import TrendReq
        keywords = keywords[:5]
        pt = TrendReq(hl="zh-TW", tz=480, timeout=(5, 15))
        pt.build_payload(keywords, timeframe=timeframe, geo=geo)
        df = pt.interest_over_time()
        if df.empty:
            return f"找不到「{', '.join(keywords)}」的趨勢資料"
        lines = [f"📈 Google Trends（{geo}，{timeframe}）\n"]
        for kw in keywords:
            if kw not in df.columns:
                continue
            avg = df[kw].mean()
            peak = df[kw].max()
            peak_date = str(df[kw].idxmax())[:10]
            recent = df[kw].iloc[-4:].mean()
            trend = "上升 📈" if df[kw].iloc[-1] > df[kw].iloc[-8] else "下降 📉"
            lines.append(
                f"🔍 {kw}\n"
                f"   平均熱度：{avg:.0f}　峰值：{peak}（{peak_date}）\n"
                f"   近期熱度：{recent:.0f}　趨勢：{trend}"
            )
        # 相關搜尋
        try:
            related = pt.related_queries()
            for kw in keywords[:2]:
                if kw in related and related[kw]["top"] is not None:
                    top_q = related[kw]["top"]["query"].head(3).tolist()
                    lines.append(f"\n「{kw}」相關搜尋：{', '.join(top_q)}")
        except Exception:
            pass
        return "\n".join(lines)
    except Exception as e:
        return f"Google Trends 查詢失敗：{e}"


def fetch_health_research(topic: str, lang: str = "zh-tw") -> str:
    try:
        from ddgs import DDGS
        region = "tw-tzh" if lang == "zh-tw" else "us-en"
        queries = [
            f"{topic} 症狀 原因 說明",
            f"{topic} 治療 建議 注意事項",
            f"{topic} 衛福部 醫學 研究",
        ]
        all_hits = []
        with DDGS() as ddgs:
            for q in queries:
                hits = list(ddgs.text(q, region=region, max_results=3))
                all_hits.extend(hits)
        lines = [
            f"🏥 健康資訊：{topic}\n",
            f"⚠️ 以下資訊僅供參考，不替代醫師診斷，如有不適請就醫。\n",
        ]
        seen = set()
        for h in all_hits[:6]:
            title = h.get("title", "")
            if title in seen:
                continue
            seen.add(title)
            body = h.get("body", "")[:250]
            lines.append(f"• {title}\n  {body}\n")
        return "\n".join(lines)
    except Exception as e:
        return f"健康資訊搜尋失敗：{e}"


def fetch_image(prompt: str, width: int = 512, height: int = 512):
    hf_token = os.getenv("HF_TOKEN")
    headers = {"Authorization": f"Bearer {hf_token}"}
    payload = {"inputs": prompt}
    models = [
        "black-forest-labs/FLUX.1-schnell",
        "stabilityai/stable-diffusion-xl-base-1.0",
    ]
    for model in models:
        for attempt in range(2):
            try:
                res = requests.post(
                    f"https://router.huggingface.co/hf-inference/models/{model}",
                    headers=headers,
                    json=payload,
                    timeout=120
                )
                if res.status_code == 200 and res.headers.get("content-type", "").startswith("image"):
                    return res.content
                if res.status_code == 503:
                    time.sleep(10)
                    continue
            except Exception:
                pass
    return None


def fetch_impact_analysis(event: str, scope: list = None, lang: str = "zh-tw") -> str:
    try:
        from ddgs import DDGS
        region = "tw-tzh" if lang == "zh-tw" else "us-en"
        if not scope:
            scope = ["個人", "企業", "社會", "經濟"]
        lines = [f"🌐 影響力分析：{event}\n"]
        with DDGS() as ddgs:
            for s in scope:
                hits = list(ddgs.text(f"{event} 對{s}的影響", region=region, max_results=2))
                if hits:
                    lines.append(f"── 對{s}的影響 ──")
                    for h in hits:
                        lines.append(f"• {h.get('title','')}：{h.get('body','')[:200]}")
                    lines.append("")
        return "\n".join(lines)
    except Exception as e:
        return f"影響力分析失敗：{e}"


def fetch_inflation_adjusted(nominal_return: float, years: int, amount: float,
                              inflation_rate: float = 2.0) -> str:
    try:
        # 費雪方程式：實質報酬 ≈ 名目報酬 - 通膨率
        real_return = ((1 + nominal_return / 100) / (1 + inflation_rate / 100) - 1) * 100
        # 名目終值
        nominal_fv = amount * (1 + nominal_return / 100) ** years
        # 實質終值（通膨調整後的購買力）
        real_fv = amount * (1 + real_return / 100) ** years
        # 通膨吃掉的部分
        inflation_loss = nominal_fv - real_fv
        lines = [
            f"📉 通膨調整報酬試算\n",
            f"本金：{amount:,.0f} 元",
            f"名目報酬率：{nominal_return}%　通膨率：{inflation_rate}%",
            f"實質報酬率：{real_return:.2f}%　期間：{years} 年",
            f"",
            f"── {years}年後 ──",
            f"名目終值：{nominal_fv:,.0f} 元",
            f"實質購買力：{real_fv:,.0f} 元",
            f"通膨吃掉：{inflation_loss:,.0f} 元（{inflation_loss/nominal_fv*100:.1f}%）",
            f"",
            f"💡 今天 {amount:,.0f} 元的東西，{years}年後需要 {amount*(1+inflation_rate/100)**years:,.0f} 元才買得到",
        ]
        return "\n".join(lines)
    except Exception as e:
        return f"通膨調整試算失敗：{e}"


def fetch_institutional(symbol: str = "", date: str = "") -> str:
    """台股三大法人買賣超"""
    try:
        import datetime
        if not date:
            date = dt.date.today().strftime("%Y%m%d")
        headers = {"User-Agent": "Mozilla/5.0"}

        if symbol:
            # 個股三大法人
            url = f"https://www.twse.com.tw/rwd/zh/fund/T86?response=json&date={date}&selectType=ALL"
            resp = requests.get(url, timeout=10, headers=headers)
            data = resp.json()
            if data.get("stat") != "OK":
                return f"查無資料（{date}，可能為非交易日）"
            rows = data.get("data", [])
            target = None
            for row in rows:
                if str(row[0]).strip() == str(symbol).strip():
                    target = row
                    break
            if not target:
                return f"找不到 {symbol} 的三大法人資料"
            foreign = int(target[4].replace(",", "").replace("+", ""))
            trust   = int(target[10].replace(",", "").replace("+", ""))
            dealer  = int(target[13].replace(",", "").replace("+", "")) if len(target) > 13 else 0
            total   = foreign + trust + dealer
            arrow = lambda v: "▲" if v >= 0 else "▼"
            return (
                f"📊 {symbol} 三大法人（{date[:4]}/{date[4:6]}/{date[6:]}）\n"
                f"外資：{arrow(foreign)} {abs(foreign):,} 張\n"
                f"投信：{arrow(trust)} {abs(trust):,} 張\n"
                f"自營：{arrow(dealer)} {abs(dealer):,} 張\n"
                f"合計：{arrow(total)} {abs(total):,} 張"
            )
        else:
            # 市場整體
            url = f"https://www.twse.com.tw/rwd/zh/fund/BFI82U?response=json&type=day&dayDate={date}"
            resp = requests.get(url, timeout=10, headers=headers)
            data = resp.json()
            if data.get("stat") != "OK":
                return f"查無三大法人整體資料（{date}，可能為非交易日）"
            rows = data.get("data", [])
            lines = [f"📊 台股三大法人整體買賣超（{date[:4]}/{date[4:6]}/{date[6:]}）\n"]
            for row in rows:
                name = row[0]
                buy  = row[1].replace(",", "")
                sell = row[2].replace(",", "")
                diff = row[3].replace(",", "")
                val  = int(diff.replace("+", "")) if diff.replace("+","").replace("-","").isdigit() else 0
                arrow = "▲" if val >= 0 else "▼"
                lines.append(f"{name}：{arrow} {diff} 元")
            return "\n".join(lines)
    except Exception as e:
        return f"三大法人查詢失敗：{e}"


def fetch_ipo(count: int = 10) -> str:
    """近期 IPO 行事曆"""
    try:
        import feedparser
        count = min(count, 20)
        results = []

        # 用 Google 新聞搜尋 IPO 資訊
        url = f"https://news.google.com/rss/search?q=IPO+新股+上市&hl=zh-Hant&gl=TW&ceid=TW:zh-Hant"
        feed = feedparser.parse(url)
        lines = ["🆕 近期 IPO / 新股資訊\n"]
        for entry in feed.entries[:count]:
            title = entry.get("title", "").split(" - ")[0]
            pub   = entry.get("published", "")[:16]
            lines.append(f"• {title}（{pub}）")

        # 補充美股 IPO（用 DuckDuckGo）
        try:
            from ddgs import DDGS
            with DDGS() as ddgs:
                for r in ddgs.text("upcoming IPO 2026 stock market", region="us-en", max_results=5):
                    title = r.get("title", "")
                    body  = r.get("body", "")[:100]
                    lines.append(f"🇺🇸 {title}\n   {body}")
        except Exception:
            pass

        return "\n".join(lines) if len(lines) > 1 else "暫無 IPO 資訊"
    except Exception as e:
        return f"IPO 查詢失敗：{e}"


def fetch_job_market(job_title: str, location: str = "台灣", lang: str = "zh-tw") -> str:
    try:
        from ddgs import DDGS
        region = "tw-tzh" if lang == "zh-tw" else "us-en"
        queries = {
            "薪資行情": f"{location} {job_title} 薪資 薪水 行情",
            "技能需求": f"{job_title} 必備技能 技術要求 條件",
            "市場需求": f"{location} {job_title} 職缺 需求 前景",
            "未來趨勢": f"{job_title} 產業趨勢 未來發展 AI影響",
        }
        lines = [f"💼 職涯市場分析：{job_title}（{location}）\n"]
        with DDGS() as ddgs:
            for sec, q in queries.items():
                hits = list(ddgs.text(q, region=region, max_results=2))
                if hits:
                    lines.append(f"── {sec} ──")
                    for h in hits:
                        lines.append(f"• {h.get('title','')}：{h.get('body','')[:200]}")
                    lines.append("")
        return "\n".join(lines)
    except Exception as e:
        return f"職涯市場分析失敗：{e}"


def fetch_key_insights(topic: str, count: int = 5, lang: str = "zh-tw") -> str:
    try:
        from ddgs import DDGS
        region = "tw-tzh" if lang == "zh-tw" else "us-en"
        queries = [
            f"{topic} 關鍵發現 重要結論",
            f"{topic} 研究結果 數據 證明",
            f"{topic} 專家觀點 核心",
        ]
        all_hits = []
        with DDGS() as ddgs:
            for q in queries:
                hits = list(ddgs.text(q, region=region, max_results=4))
                all_hits.extend(hits)
        # 找含數字或強烈結論的句子（通常是洞察）
        import re
        insight_kw = ["發現", "證明", "研究", "顯示", "指出", "表明", "關鍵", "重要", "核心",
                      "shows", "reveals", "found", "key", "critical", "significant"]
        candidates = []
        for h in all_hits:
            body = h.get("body", "")
            title = h.get("title", "")
            score = sum(body.count(k) + title.count(k) for k in insight_kw)
            has_num = bool(re.search(r'\d+[%倍億萬]', body))
            candidates.append((score + (2 if has_num else 0), title, body[:200]))
        candidates.sort(reverse=True)
        lines = [f"💡 核心洞察：{topic}\n"]
        seen_t = set()
        idx = 1
        for score, title, body in candidates:
            if title in seen_t or not body:
                continue
            seen_t.add(title)
            lines.append(f"#{idx} {title}\n   → {body}\n")
            idx += 1
            if idx > count:
                break
        lines.append(f"洞察來自 {len(all_hits)} 筆資料，依相關性與數據強度排序")
        return "\n".join(lines)
    except Exception as e:
        return f"洞察萃取失敗：{e}"


def fetch_law_research(topic: str, jurisdiction: str = "台灣", lang: str = "zh-tw") -> str:
    try:
        from ddgs import DDGS
        region = "tw-tzh" if lang == "zh-tw" else "us-en"
        queries = [
            f"{jurisdiction} {topic} 法律 法規 條文",
            f"{jurisdiction} {topic} 判例 實務 見解",
            f"{topic} 法律問題 解答",
        ]
        all_hits = []
        with DDGS() as ddgs:
            for q in queries:
                hits = list(ddgs.text(q, region=region, max_results=3))
                all_hits.extend(hits)
        lines = [
            f"⚖️ 法規查詢：{topic}（{jurisdiction}）\n",
            f"⚠️ 以下資訊僅供參考，不構成法律意見，具體情況建議諮詢律師。\n",
        ]
        seen = set()
        for h in all_hits[:6]:
            title = h.get("title", "")
            if title in seen:
                continue
            seen.add(title)
            body = h.get("body", "")[:250]
            lines.append(f"• {title}\n  {body}\n")
        return "\n".join(lines)
    except Exception as e:
        return f"法規查詢失敗：{e}"


def fetch_loan_calculator(principal: float, annual_rate: float, years: int,
                           loan_type: str = "等額本息") -> str:
    try:
        p = principal * 10000
        r = annual_rate / 100 / 12
        n = years * 12
        if n <= 0:
            return "⚠️ 貸款年數必須大於 0"
        total_interest = 0
        lines = [f"🏦 貸款試算（{loan_type}）\n",
                 f"貸款金額：{principal:.0f} 萬元",
                 f"年利率：{annual_rate}%　期數：{n} 期（{years} 年）\n"]
        if loan_type == "等額本息":
            if r > 0:
                payment = p * r * (1 + r) ** n / ((1 + r) ** n - 1)
            else:
                payment = p / n
            total_pay = payment * n
            total_interest = total_pay - p
            lines += [
                f"── 等額本息 ──",
                f"每月還款：{payment:,.0f} 元",
                f"總還款額：{total_pay/10000:.2f} 萬元",
                f"總利息：{total_interest/10000:.2f} 萬元",
            ]
        else:  # 等額本金
            principal_payment = p / n
            first_payment = principal_payment + p * r
            last_payment = principal_payment + principal_payment * r
            total_interest = sum(principal_payment * r * (n - i) for i in range(n))
            lines += [
                f"── 等額本金 ──",
                f"每期本金：{principal_payment:,.0f} 元",
                f"第1期還款：{first_payment:,.0f} 元",
                f"最後1期還款：{last_payment:,.0f} 元",
                f"總利息：{total_interest/10000:.2f} 萬元",
                f"總還款額：{(p + total_interest)/10000:.2f} 萬元",
            ]
        return "\n".join(lines)
    except Exception as e:
        return f"貸款試算失敗：{e}"


def fetch_macro(indicator: str) -> str:
    try:
        if indicator == "fed_rate":
            import yfinance as yf
            ticker = yf.Ticker("^IRX")
            hist = ticker.history(period="1mo")
            if not hist.empty:
                rate = hist["Close"].iloc[-1]
                prev = hist["Close"].iloc[-2] if len(hist) > 1 else rate
                change = rate - prev
                return (
                    f"🏦 美國短期利率（13週國庫券）\n"
                    f"當前：{rate:.2f}%　前日：{prev:.2f}%　變化：{change:+.2f}%\n"
                    f"近1月高低：{hist['Low'].min():.2f}% ~ {hist['High'].max():.2f}%\n"
                    f"（聯邦基金目標利率請參考 federalreserve.gov）"
                )
            return "無法取得利率資料"
        if indicator == "nonfarm":
            return (
                "📊 美國非農就業（Non-Farm Payrolls）\n"
                "每月第一個週五由美國勞工部公布。\n"
                "建議用 get_finance_news 工具搜尋最新數據，或查詢 bls.gov。"
            )
        wb_map = {
            "cpi":          ("美國 CPI 通膨年增率",  "FP.CPI.TOTL.ZG",  "US"),
            "unemployment": ("美國失業率",            "SL.UEM.TOTL.ZS",  "US"),
            "gdp":          ("美國 GDP 年增率",       "NY.GDP.MKTP.KD.ZG","US"),
        }
        label, wb_code, country = wb_map[indicator]
        url = f"https://api.worldbank.org/v2/country/{country}/indicator/{wb_code}?format=json&mrv=5"
        resp = requests.get(url, timeout=10)
        data = resp.json()
        if len(data) < 2 or not data[1]:
            return f"無法取得 {label} 資料"
        entries = [e for e in data[1] if e.get("value") is not None][:5]
        lines = [f"📈 {label}（世界銀行）"]
        for e in entries:
            lines.append(f"{e['date']}：{e['value']:.2f}%")
        return "\n".join(lines)
    except Exception as e:
        return f"總經指標查詢失敗：{e}"


def fetch_margin_trading(symbol: str, date: str = "") -> str:
    """台股融資融券餘額"""
    try:
        import datetime
        if not date:
            date = dt.date.today().strftime("%Y%m%d")
        # TWSE 融資融券 API
        url = f"https://www.twse.com.tw/rwd/zh/marginTrading/MI_MARGN?response=json&date={date}&selectType=ALL"
        resp = requests.get(url, timeout=10, headers={"User-Agent": "Mozilla/5.0"})
        data = resp.json()
        if data.get("stat") != "OK":
            return f"查無融資融券資料（{date}，可能為非交易日）"
        rows = data.get("data", [])
        target = None
        for row in rows:
            if str(row[0]).strip() == str(symbol).strip():
                target = row
                break
        if not target:
            return f"找不到 {symbol} 的融資融券資料"

        # 欄位：代號, 名稱, 融資買進, 融資賣出, 融資現金償還, 融資餘額, 融資限額,
        #        融券賣出, 融券買進, 融券現券償還, 融券餘額, 融券限額, 資券互抵
        name     = target[1]
        loan_bal = target[5].replace(",", "")   # 融資餘額（千股）
        short_bal= target[10].replace(",", "")  # 融券餘額（千股）
        loan_buy = target[2].replace(",", "")
        loan_sell= target[3].replace(",", "")
        short_sell=target[7].replace(",", "")
        short_buy= target[8].replace(",", "")

        return (
            f"📋 {name}（{symbol}）融資融券（{date[:4]}/{date[4:6]}/{date[6:]}）\n\n"
            f"── 融資（散戶多單）──\n"
            f"餘額：{loan_bal} 千股\n"
            f"今買：{loan_buy} 千股　今賣：{loan_sell} 千股\n\n"
            f"── 融券（放空）──\n"
            f"餘額：{short_bal} 千股\n"
            f"今賣：{short_sell} 千股　今買：{short_buy} 千股\n\n"
            f"資券比：{int(loan_bal)/(int(short_bal) if int(short_bal) > 0 else 1):.1f}x"
            if loan_bal.isdigit() and short_bal.isdigit() else ""
        )
    except Exception as e:
        return f"融資融券查詢失敗：{e}"


def fetch_market_sentiment() -> str:
    try:
        import yfinance as yf
        # VIX
        vix_hist = yf.Ticker("^VIX").history(period="5d")
        vix = vix_hist["Close"].iloc[-1] if not vix_hist.empty else None
        vix_prev = vix_hist["Close"].iloc[-2] if len(vix_hist) > 1 else vix

        if vix:
            if vix >= 40: vix_note = "極度恐慌 😱"
            elif vix >= 30: vix_note = "高度恐慌 😰"
            elif vix >= 20: vix_note = "輕度緊張 😟"
            elif vix >= 12: vix_note = "正常 😐"
            else: vix_note = "市場過度樂觀 😎"

        # Fear & Greed Index (CNN)
        try:
            fg_resp = requests.get(
                "https://production.dataviz.cnn.io/index/fearandgreed/graphdata",
                headers={"User-Agent": "Mozilla/5.0"}, timeout=8
            )
            fg_data = fg_resp.json()
            fg_score = fg_data["fear_and_greed"]["score"]
            fg_rating = fg_data["fear_and_greed"]["rating"]
            fg_prev = fg_data["fear_and_greed"]["previous_close"]
            fg_str = f"{fg_score:.0f}/100 — {fg_rating}"
            fg_change = fg_score - fg_prev
        except Exception:
            fg_str = "無法取得"
            fg_change = 0

        # S&P500 and Nasdaq 當日走勢
        sp_hist = yf.Ticker("^GSPC").history(period="5d")
        ndx_hist = yf.Ticker("^IXIC").history(period="5d")
        sp_chg = ((sp_hist["Close"].iloc[-1] / sp_hist["Close"].iloc[-2]) - 1) * 100 if len(sp_hist) > 1 else 0
        ndx_chg = ((ndx_hist["Close"].iloc[-1] / ndx_hist["Close"].iloc[-2]) - 1) * 100 if len(ndx_hist) > 1 else 0

        lines = ["🌡 市場情緒儀表板\n"]
        if vix:
            vix_chg = vix - vix_prev if vix_prev else 0
            lines.append(f"VIX 波動率：{vix:.2f}（{vix_note}）{vix_chg:+.2f}")
        lines.append(f"恐慌貪婪指數：{fg_str}（{fg_change:+.1f}）")
        lines.append(f"\nS&P 500：{sp_chg:+.2f}%")
        lines.append(f"Nasdaq：{ndx_chg:+.2f}%")

        return "\n".join(lines)
    except Exception as e:
        return f"市場情緒查詢失敗：{e}"


def fetch_money_flow(symbol: str) -> str:
    try:
        import yfinance as yf
        ticker = yf.Ticker(symbol)
        hist = ticker.history(period="20d")
        info = ticker.info
        name = info.get("shortName") or symbol

        if hist.empty:
            return f"找不到 {symbol} 資料"

        # 資金流向：(收盤-開盤)/當日振幅 * 成交量 估算買賣壓
        typical = (hist["High"] + hist["Low"] + hist["Close"]) / 3
        raw_mf = typical * hist["Volume"]
        pos_mf = raw_mf[hist["Close"] >= hist["Open"]].tail(10).sum()
        neg_mf = raw_mf[hist["Close"] < hist["Open"]].tail(10).sum()
        mfi = 100 - (100 / (1 + pos_mf / neg_mf)) if neg_mf else 100

        today_vol = hist["Volume"].iloc[-1]
        avg_vol = hist["Volume"].tail(20).mean()
        vol_ratio = today_vol / avg_vol if avg_vol else 1

        price = hist["Close"].iloc[-1]
        chg = (price / hist["Close"].iloc[-2] - 1) * 100 if len(hist) > 1 else 0

        flow = "淨流入 📥" if mfi > 55 else "淨流出 📤" if mfi < 45 else "中性"
        return (
            f"💰 {name}（{symbol}）資金流向\n\n"
            f"今日漲跌：{chg:+.2f}%\n"
            f"成交量：{today_vol:,}（均量 {vol_ratio:.1f}x）\n"
            f"資金流向指標（MFI）：{mfi:.1f} → {flow}\n"
            f"近10日正向資金：{pos_mf/1e8:.1f}億\n"
            f"近10日負向資金：{neg_mf/1e8:.1f}億"
        )
    except Exception as e:
        return f"資金流向查詢失敗：{e}"


def fetch_narrative_builder(topic: str, key_message: str = "", audience: str = "", lang: str = "zh-tw") -> str:
    try:
        from ddgs import DDGS
        region = "tw-tzh" if lang == "zh-tw" else "us-en"
        with DDGS() as ddgs:
            hits = list(ddgs.text(f"{topic} 問題 現況 挑戰 衝突", region=region, max_results=5))
        context = {
            "問題/現況": [],
            "衝突/張力": [],
            "洞察/轉折": [],
            "結論/行動": [],
        }
        for h in hits[:5]:
            body = h.get("body","")[:200]
            title = h.get("title","")
            context["問題/現況"].append(f"{title}：{body}")
        lines = [
            f"📖 敘事架構：{topic}\n",
        ]
        if key_message:
            lines.append(f"核心訊息：{key_message}")
        if audience:
            lines.append(f"目標受眾：{audience}")
        lines.append("")
        lines += [
            f"══ 第一幕：問題／現況 ══",
            f"（建立共鳴，讓受眾認識到問題的存在）",
        ]
        for item in context["問題/現況"][:2]:
            lines.append(f"• {item[:180]}")
        lines += [
            f"\n══ 第二幕：衝突／張力 ══",
            f"（說明為何現有方案不夠，製造戲劇張力）",
            f"• 現有做法的局限：需進一步搜尋或分析",
            f"\n══ 第三幕：洞察／轉折 ══",
            f"（提出新視角或解法，這是故事的核心）",
            f"• 核心洞察：{key_message or '待Claude結合數據補充'}",
            f"\n══ 第四幕：結論／行動 ══",
            f"（清楚的號召行動或結語）",
            f"• 建議行動：根據以上分析，Claude 將在回應中提出具體建議",
        ]
        return "\n".join(lines)
    except Exception as e:
        return f"敘事架構失敗：{e}"


def fetch_ocr_click(target_text: str, monitor: int = 1, click_type: str = "click", region: list = None) -> str:
    try:
        import pyautogui
        import pytesseract
        from PIL import Image
        import numpy as np

        # 截圖
        mon_map = {1: 0, 2: 1, 3: 2}
        try:
            import dxcam
            cam = dxcam.create(device_idx=mon_map.get(monitor, 0))
            frame = cam.grab()
            cam.release()
            if frame is None:
                raise Exception("dxcam grab failed")
            img = Image.fromarray(frame)
        except Exception:
            import mss
            with mss.mss() as sct:
                monitors = sct.monitors
                mon = monitors[monitor] if monitor < len(monitors) else monitors[1]
                if region:
                    mon = {"left": mon["left"] + region[0], "top": mon["top"] + region[1],
                           "width": region[2], "height": region[3]}
                shot = sct.grab(mon)
                img = Image.frombytes("RGB", shot.size, shot.bgra, "raw", "BGRX")

        # OCR
        data = pytesseract.image_to_data(img, lang="chi_tra+eng", output_type=pytesseract.Output.DICT)
        found = []
        for i, text in enumerate(data["text"]):
            if target_text.lower() in text.lower() and data["conf"][i] > 30:
                x = data["left"][i] + data["width"][i] // 2
                y = data["top"][i] + data["height"][i] // 2
                found.append((x, y, data["conf"][i], text))

        if not found:
            return f"OCR找不到文字「{target_text}」，請確認文字正確或改用 vision_locate"

        # 取信心最高的結果，換算為螢幕絕對座標
        best = max(found, key=lambda f: f[2])
        cx, cy = best[0], best[1]

        # 加上螢幕偏移
        try:
            import mss
            with mss.mss() as sct:
                mon = sct.monitors[monitor] if monitor < len(sct.monitors) else sct.monitors[1]
                abs_x = mon["left"] + cx
                abs_y = mon["top"] + cy
        except Exception:
            abs_x, abs_y = cx, cy

        # 執行點擊（SendInput，支援螢幕2負座標）
        _si_universal(abs_x, abs_y, click_type)
        return f"✅ OCR找到「{best[3]}」（信心{best[2]}%），已在 ({abs_x}, {abs_y}) 執行 {click_type}"
    except Exception as e:
        return f"OCR點擊失敗：{e}"


def fetch_opinion_writer(topic: str, stance: str = "中立", style: str = "正式") -> str:
    try:
        from ddgs import DDGS
        hits = []
        with DDGS() as ddgs:
            hits = list(ddgs.text(f"{topic} 分析 觀點 評論", region="tw-tzh", max_results=6))
        context = "\n".join(f"- {h.get('title','')}：{h.get('body','')[:200]}" for h in hits)
        lines = [f"✍️ 觀點撰寫：{topic}",
                 f"立場：{stance}　文風：{style}\n",
                 f"── 資料基礎 ──",
                 context,
                 f"\n── {stance}立場分析 ──",
        ]
        if stance == "支持":
            lines.append(f"從現有資料來看，「{topic}」有其正面價值。以下論點支持此立場：")
        elif stance == "反對":
            lines.append(f"從現有資料來看，「{topic}」存在值得警惕的問題。以下論點提出質疑：")
        elif stance == "批判":
            lines.append(f"以批判性視角審視「{topic}」，可發現以下值得深究之處：")
        else:
            lines.append(f"綜合現有資料，「{topic}」可從多角度理解：")
        lines.append(f"（本節由 Claude 依蒐集資料整合後發表看法）")
        return "\n".join(lines)
    except Exception as e:
        return f"觀點撰寫失敗：{e}"


def fetch_options(symbol: str, expiry: str = "") -> str:
    """選擇權鏈資料"""
    try:
        import yfinance as yf
        ticker = yf.Ticker(symbol)
        expirations = ticker.options
        if not expirations:
            return f"{symbol} 無選擇權資料"

        exp = expiry if expiry in expirations else expirations[0]
        chain = ticker.option_chain(exp)
        calls = chain.calls
        puts  = chain.puts

        price = ticker.info.get("currentPrice") or ticker.info.get("regularMarketPrice", 0)

        # 取最接近現價的 5 個履約價
        calls = calls.iloc[(calls["strike"] - price).abs().argsort()[:5]].sort_values("strike")
        puts  = puts.iloc[(puts["strike"]  - price).abs().argsort()[:5]].sort_values("strike")

        lines = [f"📈 {symbol} 選擇權（到期：{exp}，現價：{price:.2f}）\n"]
        lines.append("── Call（買權）──")
        for _, row in calls.iterrows():
            iv = f"IV {row.get('impliedVolatility', 0)*100:.0f}%" if row.get("impliedVolatility") else ""
            oi = f"OI {row.get('openInterest', 0):,}" if row.get("openInterest") else ""
            lines.append(f"  履約 {row['strike']:.0f}：{row.get('lastPrice', 0):.2f}  {iv}  {oi}")

        lines.append("\n── Put（賣權）──")
        for _, row in puts.iterrows():
            iv = f"IV {row.get('impliedVolatility', 0)*100:.0f}%" if row.get("impliedVolatility") else ""
            oi = f"OI {row.get('openInterest', 0):,}" if row.get("openInterest") else ""
            lines.append(f"  履約 {row['strike']:.0f}：{row.get('lastPrice', 0):.2f}  {iv}  {oi}")

        return "\n".join(lines)
    except Exception as e:
        return f"選擇權查詢失敗：{e}"


def fetch_person_research(name: str, context: str = "", lang: str = "zh-tw") -> str:
    try:
        from ddgs import DDGS
        region = "tw-tzh" if lang == "zh-tw" else "us-en"
        ctx = f" {context}" if context else ""
        queries = [
            f"{name}{ctx} 背景 經歷 介紹",
            f"{name}{ctx} 成就 評價 貢獻",
            f"{name}{ctx} 爭議 批評 問題",
        ]
        sections = {"背景與經歷": [], "成就與評價": [], "爭議與批評": []}
        sec_keys = list(sections.keys())
        with DDGS() as ddgs:
            for i, q in enumerate(queries):
                hits = list(ddgs.text(q, region=region, max_results=3))
                sections[sec_keys[i]] = hits
        lines = [f"👤 人物研究：{name}\n"]
        if context:
            lines.append(f"背景：{context}\n")
        for sec, hits in sections.items():
            if hits:
                lines.append(f"── {sec} ──")
                for h in hits[:2]:
                    lines.append(f"• {h.get('title','')}：{h.get('body','')[:200]}")
                lines.append("")
        return "\n".join(lines)
    except Exception as e:
        return f"人物研究失敗：{e}"


def fetch_portfolio_risk(holdings: list, period: str = "1y") -> str:
    try:
        import yfinance as yf
        import numpy as np
        import pandas as pd

        if not holdings or not isinstance(holdings, list):
            return "⚠️ 請提供持股清單，格式：[{\"symbol\": \"AAPL\", \"weight\": 50}, ...]"
        if isinstance(holdings[0], str):
            holdings = [{"symbol": s, "weight": 1} for s in holdings]
        weights = [h.get("weight", 1) if isinstance(h, dict) else 1 for h in holdings]
        symbols = [h.get("symbol", h) if isinstance(h, dict) else str(h) for h in holdings]

        # 標準化權重
        total_w = sum(weights)
        if total_w == 0:
            return "⚠️ 所有持股權重為 0，請設定有效權重"
        weights = [w / total_w for w in weights]

        data = {}
        for sym in symbols:
            hist = yf.Ticker(sym).history(period=period)
            if not hist.empty:
                data[sym] = hist["Close"].pct_change().dropna()

        if not data:
            return "無法取得任何股票資料"

        df = pd.DataFrame(data).dropna()
        valid_syms = list(df.columns)
        valid_weights = [weights[symbols.index(s)] for s in valid_syms]
        valid_weights = [w / sum(valid_weights) for w in valid_weights]

        # 組合報酬
        portfolio_ret = (df * valid_weights).sum(axis=1)
        annual_ret = portfolio_ret.mean() * 252 * 100
        annual_vol = portfolio_ret.std() * (252 ** 0.5) * 100
        sharpe = (annual_ret/100 - 0.05) / (annual_vol/100) if annual_vol else 0
        max_dd = ((1 + portfolio_ret).cumprod() / (1 + portfolio_ret).cumprod().cummax() - 1).min() * 100
        var_95 = np.percentile(portfolio_ret, 5) * 100

        # 相關性
        corr = df.corr()

        lines = [f"📊 投資組合風險分析（{period}）\n"]
        lines.append("持倉配置：")
        for sym, w in zip(valid_syms, valid_weights):
            lines.append(f"  {sym}：{w*100:.1f}%")

        lines += [
            f"\n── 組合風險指標 ──",
            f"年化報酬：{annual_ret:+.2f}%",
            f"年化波動率：{annual_vol:.2f}%",
            f"夏普比率：{sharpe:.2f}",
            f"最大回撤：{max_dd:.2f}%",
            f"VaR 95%（單日）：{var_95:.2f}%",
        ]

        if len(valid_syms) >= 2:
            lines.append("\n── 相關性（越低越分散）──")
            for i, s1 in enumerate(valid_syms):
                for s2 in valid_syms[i+1:]:
                    v = corr.loc[s1, s2]
                    note = "⚠️高度相關" if v > 0.7 else "✅低相關" if v < 0.3 else ""
                    lines.append(f"  {s1} & {s2}：{v:.3f} {note}")

        return "\n".join(lines)
    except Exception as e:
        return f"投資組合風險分析失敗：{e}"


def fetch_position_statement(issue: str, stance: str, lang: str = "zh-tw") -> str:
    try:
        from ddgs import DDGS
        region = "tw-tzh" if lang == "zh-tw" else "us-en"
        stance_q = "支持 贊成 優點" if stance == "支持" else ("反對 問題 缺點" if stance == "反對" else "條件 但書 前提")
        queries = [
            f"{issue} {stance_q} 論據 證據",
            f"{issue} 數據 研究 案例",
            f"{issue} 反對方 質疑 反駁",
        ]
        evidence, data, counter = [], [], []
        with DDGS() as ddgs:
            evidence = list(ddgs.text(queries[0], region=region, max_results=4))
            data = list(ddgs.text(queries[1], region=region, max_results=3))
            counter = list(ddgs.text(queries[2], region=region, max_results=3))
        lines = [
            f"📣 立場聲明：{issue}",
            f"立場：{stance}\n",
            f"══ 論點（Claim）══",
            f"對於「{issue}」，我的立場是【{stance}】，理由如下：\n",
            f"── 論據與證據 ──",
        ]
        for h in evidence[:3]:
            lines.append(f"• {h.get('title','')}：{h.get('body','')[:180]}")
        lines.append("\n── 數據支撐 ──")
        for h in data[:2]:
            lines.append(f"• {h.get('title','')}：{h.get('body','')[:180]}")
        lines.append("\n── 預判反駁與回應 ──")
        for h in counter[:2]:
            lines.append(f"反方：{h.get('title','')}：{h.get('body','')[:150]}")
        lines += [
            "\n── 結論 ──",
            f"綜合以上，{stance}「{issue}」的立場是有根據且可辯護的。",
            "（Claude 將在回應中整合以上資料，給出完整系統性論述）",
        ]
        return "\n".join(lines)
    except Exception as e:
        return f"立場聲明失敗：{e}"


def fetch_product_review(product: str, category: str = "", lang: str = "zh-tw") -> str:
    try:
        from ddgs import DDGS
        region = "tw-tzh" if lang == "zh-tw" else "us-en"
        cat = f" {category}" if category else ""
        queries = [
            f"{product}{cat} 評測 開箱 評價",
            f"{product}{cat} 優點 缺點 推薦",
            f"{product}{cat} 使用心得 評分",
        ]
        all_hits = []
        with DDGS() as ddgs:
            for q in queries:
                hits = list(ddgs.text(q, region=region, max_results=3))
                all_hits.extend(hits)
        pos_kw = ["推薦", "好用", "優秀", "值得", "滿意", "excellent", "great", "recommend"]
        neg_kw = ["不推", "缺點", "問題", "失望", "差", "poor", "bad", "issue"]
        all_text = " ".join(h.get("body","") for h in all_hits).lower()
        pos_score = sum(all_text.count(k) for k in pos_kw)
        neg_score = sum(all_text.count(k) for k in neg_kw)
        total = pos_score + neg_score
        rating = round(pos_score / total * 5, 1) if total > 0 else 3.0
        lines = [
            f"⭐ 產品評測：{product}",
            f"綜合評分：{rating}/5.0　（正面{pos_score}則 / 負面{neg_score}則）\n",
            f"── 評測彙整 ──",
        ]
        seen = set()
        for h in all_hits[:6]:
            title = h.get("title","")
            if title in seen:
                continue
            seen.add(title)
            body = h.get("body","")[:200]
            lines.append(f"• {title}：{body}\n")
        return "\n".join(lines)
    except Exception as e:
        return f"產品評測失敗：{e}"


def fetch_pros_cons_analysis(subject: str, context: str = "", lang: str = "zh-tw") -> str:
    try:
        from ddgs import DDGS
        region = "tw-tzh" if lang == "zh-tw" else "us-en"
        ctx = f" {context}" if context else ""
        pros_hits, cons_hits = [], []
        with DDGS() as ddgs:
            pros_hits = list(ddgs.text(f"{subject}{ctx} 優點 好處 支持", region=region, max_results=4))
            cons_hits = list(ddgs.text(f"{subject}{ctx} 缺點 壞處 風險 問題", region=region, max_results=4))
        lines = [f"📋 優缺點分析：{subject}\n"]
        if context:
            lines.append(f"背景：{context}\n")
        lines.append("── 優點 / 支持論點 ──")
        for h in pros_hits[:3]:
            lines.append(f"✅ {h.get('title','')}：{h.get('body','')[:150]}")
        lines.append("\n── 缺點 / 反對論點 ──")
        for h in cons_hits[:3]:
            lines.append(f"⚠️ {h.get('title','')}：{h.get('body','')[:150]}")
        confidence = "高" if len(pros_hits) + len(cons_hits) >= 6 else "中"
        lines.append(f"\n資料信心度：{confidence}（共 {len(pros_hits)+len(cons_hits)} 筆）")
        return "\n".join(lines)
    except Exception as e:
        return f"優缺點分析失敗：{e}"


def fetch_read_screen(question: str = "描述螢幕上有什麼", monitor: int = 1) -> str:
    """截圖 → OCR + Vision → 回傳螢幕內容描述"""
    try:
        import anthropic, base64, io
        from PIL import Image as _PI
        img, _, _ = _cap_monitor_logical(monitor)
        if img.width > 2048:
            r = 2048 / img.width
            img = img.resize((2048, int(img.height*r)), _PI.LANCZOS)
        buf = io.BytesIO(); img.save(buf, format="JPEG", quality=92)
        if buf.tell() > 4 * 1024 * 1024:
            buf = io.BytesIO(); img.save(buf, format="JPEG", quality=80)
        b64 = base64.standard_b64encode(buf.getvalue()).decode()
        # OCR 輔助
        ocr_hint = ""
        try:
            import pytesseract
            ocr_text = pytesseract.image_to_string(img, lang="chi_tra+eng").strip()
            if ocr_text:
                ocr_hint = f"\n\nOCR偵測到的文字：{ocr_text[:500]}"
        except Exception:
            pass
        resp = anthropic.Anthropic().messages.create(
            model="claude-sonnet-4-6", max_tokens=1024,
            messages=[{"role": "user", "content": [
                {"type": "image", "source": {"type": "base64", "media_type": "image/jpeg", "data": b64}},
                {"type": "text", "text": f"這是螢幕{monitor}的截圖（{img.width}x{img.height}px）。{question}。請用繁體中文詳細描述畫面內容。{ocr_hint}"}
            ]}]
        )
        return f"📺 螢幕{monitor}內容：\n{resp.content[0].text}"
    except Exception as e:
        return f"讀取螢幕失敗：{e}"


def fetch_reits(symbol: str) -> str:
    try:
        import yfinance as yf
        ticker = yf.Ticker(symbol)
        info = ticker.info
        hist = ticker.history(period="1y")
        if hist.empty:
            return f"找不到 {symbol} 的REITs資料"
        price = hist["Close"].iloc[-1]
        price_1y = hist["Close"].iloc[0]
        ret_1y = (price / price_1y - 1) * 100
        name = info.get("longName") or info.get("shortName") or symbol
        div_yield = info.get("dividendYield") or 0
        trailing_annual_div = info.get("trailingAnnualDividendYield") or div_yield
        market_cap = info.get("marketCap") or 0
        sector = info.get("sector") or info.get("category") or "不動產"
        nav = info.get("bookValue") or info.get("navPrice") or 0
        lines = [
            f"🏢 REITs查詢：{name}\n",
            f"代號：{symbol}　類別：{sector}",
            f"現價：{price:.2f}",
        ]
        if div_yield:
            lines.append(f"股息殖利率：{trailing_annual_div*100:.2f}%")
        if market_cap:
            lines.append(f"市值：{market_cap/1e8:.1f} 億")
        if nav:
            premium = (price / nav - 1) * 100
            lines.append(f"NAV：{nav:.2f}（折溢價：{premium:+.1f}%）")
        lines += [
            f"",
            f"近1年報酬：{ret_1y:+.2f}%",
        ]
        return "\n".join(lines)
    except Exception as e:
        return f"REITs查詢失敗：{e}"


def fetch_research_report(topic: str, purpose: str = "一般研究", lang: str = "zh-tw") -> str:
    try:
        from ddgs import DDGS
        region = "tw-tzh" if lang == "zh-tw" else "us-en"
        sections = {
            "背景與定義": f"{topic} 定義 介紹 背景",
            "現況與數據": f"{topic} 現況 統計 數據 規模",
            "主要發現": f"{topic} 研究 發現 結果 報告",
            "爭議與挑戰": f"{topic} 問題 挑戰 爭議",
            "趨勢與展望": f"{topic} 趨勢 未來 預測 展望",
        }
        collected = {}
        with DDGS() as ddgs:
            for sec, q in sections.items():
                hits = list(ddgs.text(q, region=region, max_results=3))
                collected[sec] = hits
        lines = [
            f"📄 研究報告：{topic}",
            f"目的：{purpose}\n",
            f"═══ 執行摘要 ═══",
            f"本報告針對「{topic}」進行多面向資料蒐集，涵蓋背景、數據、爭議與展望。\n",
        ]
        for sec, hits in collected.items():
            lines.append(f"── {sec} ──")
            for h in hits[:2]:
                lines.append(f"• {h.get('title','')}：{h.get('body','')[:200]}")
            lines.append("")
        lines += [
            "═══ 結論與建議 ═══",
            f"根據蒐集資料，「{topic}」是一個值得深入關注的議題。",
            f"建議進一步參閱原始來源以獲得更完整資訊。",
        ]
        return "\n".join(lines)
    except Exception as e:
        return f"研究報告生成失敗：{e}"


def fetch_retirement_calculator(current_age: int, current_savings: float, monthly_save: float,
                                 retire_age: int = 65, annual_return: float = 6.0,
                                 monthly_expense: float = 50000) -> str:
    try:
        years = retire_age - current_age
        if years <= 0:
            return "退休年齡必須大於目前年齡"
        r_monthly = annual_return / 100 / 12
        # 現有資產複利成長
        future_current = current_savings * 10000 * ((1 + r_monthly) ** (years * 12))
        # 每月儲蓄複利成長（年金終值）
        if r_monthly > 0:
            future_monthly = monthly_save * (((1 + r_monthly) ** (years * 12) - 1) / r_monthly)
        else:
            future_monthly = monthly_save * years * 12
        total_at_retire = future_current + future_monthly
        # 退休後可用年數（假設活到85歲）
        retire_years = 85 - retire_age
        total_needed = monthly_expense * 12 * retire_years
        gap = total_at_retire - total_needed
        status = "✅ 達標" if gap >= 0 else "⚠️ 不足"

        lines = [
            f"🏖️ 退休規劃試算\n",
            f"目前年齡：{current_age} 歲　預計退休：{retire_age} 歲",
            f"距退休：{years} 年　預期報酬：{annual_return}%/年",
            f"",
            f"── 退休時預估資產 ──",
            f"現有資產成長至：{future_current/10000:.0f} 萬元",
            f"累積儲蓄成長至：{future_monthly/10000:.0f} 萬元",
            f"退休時總資產：{total_at_retire/10000:.0f} 萬元",
            f"",
            f"── 退休所需評估（活至85歲）──",
            f"月生活費：{monthly_expense:,.0f} 元",
            f"退休後需要：{total_needed/10000:.0f} 萬元",
            f"缺口/剩餘：{gap/10000:+.0f} 萬元　{status}",
        ]
        if gap < 0:
            extra = abs(gap) / (years * 12)
            lines.append(f"\n每月需額外存：{extra:,.0f} 元 才能達標")
        return "\n".join(lines)
    except Exception as e:
        return f"退休試算失敗：{e}"


def fetch_risk_metrics(symbol: str, period: str = "1y") -> str:
    try:
        import yfinance as yf
        import numpy as np

        ticker = yf.Ticker(symbol)
        hist = ticker.history(period=period)
        bench = yf.Ticker("^GSPC").history(period=period)
        name = ticker.info.get("shortName") or symbol

        if hist.empty or bench.empty:
            return f"找不到 {symbol} 資料"

        ret = hist["Close"].pct_change().dropna()
        bench_ret = bench["Close"].pct_change().dropna()
        aligned = ret.align(bench_ret, join="inner")
        ret, bench_ret = aligned[0], aligned[1]

        # Beta
        cov = np.cov(ret, bench_ret)
        beta = cov[0][1] / cov[1][1] if cov[1][1] != 0 else 0

        # 年化報酬 & 波動率
        annual_ret = ret.mean() * 252 * 100
        annual_vol = ret.std() * (252 ** 0.5) * 100

        # 夏普比率（無風險利率 5%）
        rf = 0.05
        sharpe = (annual_ret/100 - rf) / (annual_vol/100) if annual_vol != 0 else 0

        # 最大回撤
        cumret = (1 + ret).cumprod()
        max_dd = ((cumret / cumret.cummax()) - 1).min() * 100

        # VaR 95%
        var_95 = np.percentile(ret, 5) * 100

        return (
            f"⚖️ {name}（{symbol}）風險指標（{period}）\n\n"
            f"Beta（市場敏感度）：{beta:.2f}{'（高波動）' if beta>1.2 else '（低波動）' if beta<0.8 else '（接近大盤）'}\n"
            f"年化報酬：{annual_ret:+.2f}%\n"
            f"年化波動率：{annual_vol:.2f}%\n"
            f"夏普比率：{sharpe:.2f}{'（優秀）' if sharpe>1 else '（尚可）' if sharpe>0.5 else '（偏低）'}\n"
            f"最大回撤：{max_dd:.2f}%\n"
            f"VaR 95%（單日）：{var_95:.2f}%"
        )
    except Exception as e:
        return f"風險指標計算失敗：{e}"


def fetch_scenario_planning(topic: str, horizon: str = "", lang: str = "zh-tw") -> str:
    try:
        from ddgs import DDGS
        region = "tw-tzh" if lang == "zh-tw" else "us-en"
        h_tag = f" {horizon}" if horizon else ""
        scenarios = {
            "樂觀情境": f"{topic}{h_tag} 最好情況 成功 機會",
            "基準情境": f"{topic}{h_tag} 預測 可能發展 趨勢",
            "悲觀情境": f"{topic}{h_tag} 風險 失敗 最壞情況",
        }
        lines = [f"🔭 情境規劃：{topic}"]
        if horizon:
            lines.append(f"時間範圍：{horizon}")
        lines.append("")
        with DDGS() as ddgs:
            for sc, q in scenarios.items():
                hits = list(ddgs.text(q, region=region, max_results=3))
                lines.append(f"══ {sc} ══")
                for h in hits[:2]:
                    lines.append(f"• {h.get('title','')}：{h.get('body','')[:200]}")
                lines.append("")
        lines.append("💡 建議針對各情境預先制定因應策略，提高應變能力")
        return "\n".join(lines)
    except Exception as e:
        return f"情境規劃失敗：{e}"


def fetch_screen_workflow(steps: list) -> str:
    try:
        import time
        results = []
        for i, step in enumerate(steps):
            action = step.get("action", "")
            target = step.get("target", "")
            value = step.get("value", "")
            monitor = step.get("monitor", 1)
            try:
                if action == "screenshot":
                    results.append(f"步驟{i+1} screenshot：已截圖")
                elif action == "ocr_click":
                    r = fetch_ocr_click(target, monitor)
                    results.append(f"步驟{i+1} ocr_click [{target}]：{r}")
                elif action == "vision_click":
                    r = fetch_vision_locate(target, monitor, "click")
                    results.append(f"步驟{i+1} vision_click [{target}]：{r}")
                elif action == "type":
                    import pyautogui
                    pyautogui.write(value, interval=0.05)
                    results.append(f"步驟{i+1} type：已輸入「{value[:30]}」")
                elif action == "press":
                    import pyautogui
                    pyautogui.press(value)
                    results.append(f"步驟{i+1} press：已按 {value}")
                elif action == "wait":
                    secs = float(value) if value else 1.0
                    time.sleep(secs)
                    results.append(f"步驟{i+1} wait：等待 {secs}s")
                elif action == "open_app":
                    import subprocess
                    subprocess.Popen(target, shell=True)
                    time.sleep(1.5)
                    results.append(f"步驟{i+1} open_app：已開啟 {target}")
                elif action == "hotkey":
                    import pyautogui
                    keys = [k.strip() for k in value.split("+")]
                    pyautogui.hotkey(*keys)
                    results.append(f"步驟{i+1} hotkey：{value}")
                else:
                    results.append(f"步驟{i+1} 未知動作：{action}")
            except Exception as e:
                results.append(f"步驟{i+1} 失敗：{e}")
                break
        return "📋 工作流執行結果：\n" + "\n".join(results)
    except Exception as e:
        return f"螢幕工作流失敗：{e}"


def fetch_scroll_at(direction: str = "down", amount: int = 3,
                    x: int = None, y: int = None,
                    monitor: int = 1, description: str = "") -> str:
    """在指定位置滾動，支援所有螢幕包含螢幕2負座標"""
    try:
        import mss, time
        if description:
            img, ml, mt = _cap_monitor_logical(monitor)
            rx, ry = _vision_find(img, description)
            if rx is not None:
                abs_x, abs_y = ml + rx, mt + ry
            else:
                with mss.mss() as s:
                    m = s.monitors[monitor] if monitor < len(s.monitors) else s.monitors[1]
                    abs_x = m["left"] + m["width"] // 2
                    abs_y = m["top"] + m["height"] // 2
        elif x is not None and y is not None:
            with mss.mss() as s:
                m = s.monitors[monitor] if monitor < len(s.monitors) else s.monitors[1]
                abs_x = m["left"] + x; abs_y = m["top"] + y
        else:
            with mss.mss() as s:
                m = s.monitors[monitor] if monitor < len(s.monitors) else s.monitors[1]
                abs_x = m["left"] + m["width"] // 2
                abs_y = m["top"] + m["height"] // 2
        _si_scroll(abs_x, abs_y, amount, direction)
        return f"✅ 螢幕{monitor} 在({abs_x},{abs_y}) 向{direction}滾動 {amount} 格"
    except Exception as e:
        return f"滾動失敗：{e}"


def fetch_second_opinion(question: str, experts: list = None, lang: str = "zh-tw") -> str:
    try:
        from ddgs import DDGS
        region = "tw-tzh" if lang == "zh-tw" else "us-en"
        if not experts:
            experts = ["經濟學家", "心理學家", "社會學家", "科技專家", "實務工作者"]
        lines = [f"🎓 多專家視角：{question}\n"]
        with DDGS() as ddgs:
            for expert in experts[:5]:
                hits = list(ddgs.text(f"{question} {expert} 觀點 看法", region=region, max_results=2))
                lines.append(f"── {expert}的角度 ──")
                if hits:
                    for h in hits[:1]:
                        lines.append(f"• {h.get('title','')}：{h.get('body','')[:200]}")
                else:
                    lines.append(f"• 資料不足，需進一步搜尋")
                lines.append("")
        lines.append("💡 不同專業背景會產生截然不同的分析框架，綜合參考才能形成全面判斷")
        return "\n".join(lines)
    except Exception as e:
        return f"多專家視角失敗：{e}"


def fetch_sector(market: str = "us") -> str:
    """產業類股表現"""
    try:
        import yfinance as yf
        if market == "us":
            sectors = {
                "科技": "XLK", "金融": "XLF", "醫療": "XLV", "能源": "XLE",
                "工業": "XLI", "消費必需": "XLP", "消費選擇": "XLY",
                "公用事業": "XLU", "材料": "XLB", "通訊": "XLC", "房地產": "XLRE"
            }
        else:
            sectors = {
                "半導體": "00891.TW", "金融": "0055.TW", "航運": "00895.TW",
                "電動車": "00893.TW", "ESG": "00878.TW", "高息": "00919.TW",
                "科技": "0052.TW", "傳產": "0054.TW"
            }

        results = []
        for name, sym in sectors.items():
            try:
                hist = yf.Ticker(sym).history(period="2d")
                if len(hist) >= 2:
                    chg = (hist["Close"].iloc[-1] / hist["Close"].iloc[-2] - 1) * 100
                    arrow = "▲" if chg >= 0 else "▼"
                    results.append((chg, f"{arrow} {name}：{chg:+.2f}%"))
            except Exception:
                pass

        results.sort(key=lambda x: x[0], reverse=True)
        market_label = "美股" if market == "us" else "台股"
        lines = [f"🏭 {market_label}產業類股今日表現\n"]
        for _, line in results:
            lines.append(line)
        if results:
            lines.append(f"\n最強：{results[0][1].split('：')[0].strip()}")
            lines.append(f"最弱：{results[-1][1].split('：')[0].strip()}")
        return "\n".join(lines)
    except Exception as e:
        return f"類股查詢失敗：{e}"


def fetch_sentiment_scan(topic: str, lang: str = "zh-tw") -> str:
    try:
        from ddgs import DDGS
        region = "tw-tzh" if lang == "zh-tw" else "us-en"
        queries = [
            f"{topic} 正面 支持 優點",
            f"{topic} 負面 反對 批評 缺點",
            f"{topic} 民眾看法 輿論 評價",
        ]
        pos_hits, neg_hits, neutral_hits = [], [], []
        with DDGS() as ddgs:
            pos_hits = list(ddgs.text(queries[0], region=region, max_results=4))
            neg_hits = list(ddgs.text(queries[1], region=region, max_results=4))
            neutral_hits = list(ddgs.text(queries[2], region=region, max_results=3))
        total = len(pos_hits) + len(neg_hits) + len(neutral_hits)
        pos_pct = round(len(pos_hits) / total * 100) if total else 0
        neg_pct = round(len(neg_hits) / total * 100) if total else 0
        neu_pct = 100 - pos_pct - neg_pct
        lines = [
            f"📊 輿情掃描：{topic}\n",
            f"正面 {pos_pct}% ｜ 負面 {neg_pct}% ｜ 中立 {neu_pct}%\n",
            f"── 正面觀點 ──",
        ]
        for h in pos_hits[:2]:
            lines.append(f"• {h.get('title','')}：{h.get('body','')[:120]}")
        lines.append(f"\n── 負面觀點 ──")
        for h in neg_hits[:2]:
            lines.append(f"• {h.get('title','')}：{h.get('body','')[:120]}")
        lines.append(f"\n── 中立/綜合 ──")
        for h in neutral_hits[:2]:
            lines.append(f"• {h.get('title','')}：{h.get('body','')[:120]}")
        return "\n".join(lines)
    except Exception as e:
        return f"輿情掃描失敗：{e}"


def fetch_short_interest(symbol: str) -> str:
    try:
        import yfinance as yf
        info = yf.Ticker(symbol).info
        name = info.get("shortName") or symbol

        short_pct = info.get("shortPercentOfFloat", 0) or 0
        short_ratio = info.get("shortRatio", 0) or 0
        shares_short = info.get("sharesShort", 0) or 0
        shares_out = info.get("sharesOutstanding", 1) or 1

        if short_pct >= 0.20:
            risk = "極高空頭壓力（軋空機會大 ⚠️）"
        elif short_pct >= 0.10:
            risk = "高空頭比率（需留意）"
        elif short_pct >= 0.05:
            risk = "中等空頭比率"
        else:
            risk = "低空頭比率（市場偏多）"

        return (
            f"🩳 {name}（{symbol}）空頭資料\n\n"
            f"做空比率：{short_pct*100:.2f}%\n"
            f"回補天數（Short Ratio）：{short_ratio:.1f} 天\n"
            f"借券賣出股數：{shares_short:,}\n"
            f"風險評估：{risk}"
        )
    except Exception as e:
        return f"空頭資料查詢失敗：{e}"


def fetch_socratic_questioning(topic: str, depth: int = 5, lang: str = "zh-tw") -> str:
    try:
        from ddgs import DDGS
        region = "tw-tzh" if lang == "zh-tw" else "us-en"
        # 搜尋議題背景，生成有依據的問題層次
        with DDGS() as ddgs:
            hits = list(ddgs.text(f"{topic} 核心問題 爭議 本質", region=region, max_results=4))
        context = " ".join(h.get("body","")[:150] for h in hits)
        # 問題層次設計
        layers = [
            ("釐清概念", f"「{topic}」的核心定義是什麼？我們所說的究竟指的是哪個層面？"),
            ("探究假設", f"這個說法背後有哪些未被檢驗的假設？是否所有人都認同這些前提？"),
            ("檢驗證據", f"支持這個立場的證據有多可靠？有沒有相反的證據被忽略？"),
            ("探索觀點", f"從不同利益關係人的角度來看，這件事會有什麼不同的詮釋？"),
            ("追問影響", f"如果這個觀點是對的，它會帶來什麼後果？我們準備好接受這些後果了嗎？"),
            ("質疑問題本身", f"我們是否問了正確的問題？有沒有更根本的問題應該先被回答？"),
            ("尋找矛盾", f"這個立場內部有沒有自相矛盾之處？邊界條件在哪裡？"),
            ("回歸本質", f"剝除所有表象後，這個議題的最核心本質究竟是什麼？"),
        ][:depth]
        lines = [f"🏛️ 蘇格拉底式提問：{topic}\n"]
        for i, (layer_name, question) in enumerate(layers, 1):
            lines.append(f"第{i}層【{layer_name}】")
            lines.append(f"  ❓ {question}\n")
        lines.append(f"── 背景參考 ──")
        for h in hits[:2]:
            lines.append(f"• {h.get('title','')}：{h.get('body','')[:150]}")
        lines.append("\n💡 真正的理解來自不斷追問，而非接受第一個答案。")
        return "\n".join(lines)
    except Exception as e:
        return f"蘇格拉底提問失敗：{e}"


def fetch_steel_man(opposing_view: str, own_position: str = "", lang: str = "zh-tw") -> str:
    try:
        from ddgs import DDGS
        region = "tw-tzh" if lang == "zh-tw" else "us-en"
        # 搜尋支持對立觀點的最強論據
        pro_hits, counter_hits = [], []
        with DDGS() as ddgs:
            pro_hits = list(ddgs.text(f"{opposing_view} 支持 論據 最強理由", region=region, max_results=4))
            if own_position:
                counter_hits = list(ddgs.text(f"{own_position} 論據 支持", region=region, max_results=3))
        lines = [
            f"⚔️ 鋼人論證\n",
            f"對立觀點：「{opposing_view}」\n",
            f"══ 鋼人化：對方最強論點 ══",
            f"（以下為對方觀點的最有力版本，非我方立場）\n",
        ]
        for h in pro_hits[:4]:
            lines.append(f"✦ {h.get('title','')}：{h.get('body','')[:180]}\n")
        if own_position:
            lines += [
                f"══ 我方回應 ══",
                f"立場：「{own_position}」\n",
            ]
            for h in counter_hits[:3]:
                lines.append(f"→ {h.get('title','')}：{h.get('body','')[:180]}\n")
        lines.append("💡 鋼人論證要求：先真正理解對方最強版本，才能給出真正有力的回應。")
        return "\n".join(lines)
    except Exception as e:
        return f"鋼人論證失敗：{e}"


def fetch_stock(symbol: str, period: str = "1mo") -> str:
    try:
        import yfinance as yf
        ticker = yf.Ticker(symbol)
        info = ticker.info
        # 取較長歷史以便計算 RSI
        hist = ticker.history(period="3mo")

        if hist.empty:
            return f"找不到「{symbol}」的股票數據，請確認代號是否正確。"

        name = info.get("longName") or info.get("shortName") or symbol
        currency = info.get("currency", "")
        current = hist["Close"].iloc[-1]
        prev = hist["Close"].iloc[-2] if len(hist) > 1 else current
        change = current - prev
        change_pct = (change / prev) * 100 if prev else 0
        arrow = "▲" if change >= 0 else "▼"
        volume = hist["Volume"].iloc[-1]
        avg_volume = hist["Volume"].tail(20).mean()
        volume_ratio = volume / avg_volume if avg_volume else 1

        # 技術指標
        ma5 = hist["Close"].tail(5).mean()
        ma20 = hist["Close"].tail(20).mean()
        ma60 = hist["Close"].tail(60).mean() if len(hist) >= 60 else hist["Close"].mean()
        rsi = calc_rsi(hist["Close"]) if len(hist) >= 15 else None

        # 趨勢
        if ma5 > ma20 > ma60:
            trend = "強勢多頭（MA5>MA20>MA60）📈"
        elif ma5 < ma20 < ma60:
            trend = "強勢空頭（MA5<MA20<MA60）📉"
        elif ma5 > ma20:
            trend = "短線偏多（MA5>MA20）🔼"
        else:
            trend = "短線偏空（MA5<MA20）🔽"

        # RSI 解讀
        rsi_note = ""
        if rsi is not None:
            if rsi >= 80:
                rsi_note = "（嚴重超買 ⚠️）"
            elif rsi >= 70:
                rsi_note = "（超買區間）"
            elif rsi <= 20:
                rsi_note = "（嚴重超賣 💡）"
            elif rsi <= 30:
                rsi_note = "（超賣區間）"
            else:
                rsi_note = "（中性）"

        # 基本面
        market_cap = info.get("marketCap")
        pe_ratio = info.get("trailingPE")
        week52_high = info.get("fiftyTwoWeekHigh")
        week52_low = info.get("fiftyTwoWeekLow")
        high_period = hist["High"].tail(20).max()
        low_period = hist["Low"].tail(20).min()

        # 距離高低點百分比
        pct_from_high = ((current - high_period) / high_period * 100) if high_period else 0
        pct_from_low = ((current - low_period) / low_period * 100) if low_period else 0

        result = (
            f"📊 {name} ({symbol})\n"
            f"💰 現價：{current:.2f} {currency}  {arrow} {abs(change):.2f} ({change_pct:+.2f}%)\n"
            f"📦 成交量：{volume:,}（均量 {volume_ratio:.1f}x）\n"
            f"\n── 技術指標 ──\n"
            f"MA5：{ma5:.2f}　MA20：{ma20:.2f}　MA60：{ma60:.2f}\n"
            f"趨勢：{trend}\n"
        )

        if rsi is not None:
            result += f"RSI(14)：{rsi}{rsi_note}\n"

        result += (
            f"近20日高點：{high_period:.2f}（距高 {pct_from_high:.1f}%）\n"
            f"近20日低點：{low_period:.2f}（距低 +{pct_from_low:.1f}%）\n"
        )

        if week52_high and week52_low:
            result += f"52週高低：{week52_low:.2f} ~ {week52_high:.2f}\n"

        result += "\n── 基本面 ──\n"
        if market_cap:
            mc_str = f"{market_cap/1e12:.2f}兆" if market_cap >= 1e12 else f"{market_cap/1e8:.0f}億"
            result += f"市值：{mc_str} {currency}\n"
        if pe_ratio:
            result += f"本益比：{pe_ratio:.1f}\n"

        return result.strip()

    except Exception as e:
        return f"查詢「{symbol}」失敗：{str(e)}"


def fetch_stock_advanced(symbol: str, indicators: list = None) -> str:
    try:
        import yfinance as yf
        import ta as ta_lib
        if indicators is None:
            indicators = ["macd", "bb", "kd"]
        ticker = yf.Ticker(symbol)
        hist = ticker.history(period="6mo")
        if hist.empty:
            return f"找不到「{symbol}」的資料"
        name = ticker.info.get("longName") or ticker.info.get("shortName") or symbol
        close = hist["Close"]
        high = hist["High"]
        low = hist["Low"]
        current = close.iloc[-1]
        result = [f"📊 {name} ({symbol}) 進階技術分析\n現價：{current:.2f}\n"]
        if "macd" in indicators:
            macd_ind = ta_lib.trend.MACD(close)
            macd_v = macd_ind.macd().iloc[-1]
            signal_v = macd_ind.macd_signal().iloc[-1]
            hist_v = macd_ind.macd_diff().iloc[-1]
            cross = "金叉 📈" if macd_v > signal_v else "死叉 📉"
            result.append(
                f"── MACD ──\n"
                f"MACD：{macd_v:.3f}　Signal：{signal_v:.3f}　Histogram：{hist_v:.3f}\n"
                f"狀態：{cross}"
            )
        if "bb" in indicators:
            bb = ta_lib.volatility.BollingerBands(close)
            upper = bb.bollinger_hband().iloc[-1]
            mid = bb.bollinger_mavg().iloc[-1]
            lower = bb.bollinger_lband().iloc[-1]
            width = (upper - lower) / mid * 100 if mid != 0 else 0
            if current >= upper:
                bb_pos = "觸上軌（超買警示）⚠️"
            elif current <= lower:
                bb_pos = "觸下軌（超賣機會）💡"
            else:
                pos_pct = (current - lower) / (upper - lower) * 100
                bb_pos = f"通道內 {pos_pct:.0f}%"
            result.append(
                f"\n── 布林通道（BB）──\n"
                f"上軌：{upper:.2f}　中軌：{mid:.2f}　下軌：{lower:.2f}\n"
                f"帶寬：{width:.1f}%　位置：{bb_pos}"
            )
        if "kd" in indicators:
            stoch = ta_lib.momentum.StochasticOscillator(high, low, close)
            k = stoch.stoch().iloc[-1]
            d = stoch.stoch_signal().iloc[-1]
            if k > 80:
                kd_note = "超買區（K>80）⚠️"
            elif k < 20:
                kd_note = "超賣區（K<20）💡"
            else:
                kd_note = "中性區間"
            cross_kd = "K上穿D（買訊）📈" if k > d else "K下穿D（賣訊）📉"
            result.append(
                f"\n── KD 指標 ──\n"
                f"K：{k:.1f}　D：{d:.1f}\n"
                f"狀態：{kd_note}　交叉：{cross_kd}"
            )
        return "\n".join(result)
    except Exception as e:
        return f"進階技術分析失敗：{e}"


def fetch_stock_screener(criteria: str, market: str = "us") -> str:
    """選股篩選器（用 Claude 解讀條件 + yfinance 驗證）"""
    try:
        import yfinance as yf

        # 美股用 S&P500 成分股，台股用常見大型股
        if market == "us":
            # 取 S&P500 部分成分股做示範
            candidates = [
                "AAPL","MSFT","GOOGL","AMZN","NVDA","META","TSLA","BRK-B","JPM","V",
                "XOM","UNH","JNJ","WMT","MA","PG","HD","CVX","MRK","ABBV",
                "KO","PEP","BAC","PFE","AVGO","COST","TMO","MCD","ABT","CRM",
                "ACN","LIN","DHR","TXN","NEE","QCOM","PM","HON","IBM","GE",
                "ORCL","AMGN","SBUX","CAT","INTU","AMD","ISRG","NOW","MDLZ","AXP"
            ]
        else:
            candidates = [
                "2330.TW","2317.TW","2454.TW","2412.TW","2308.TW","2303.TW",
                "2881.TW","2882.TW","2886.TW","2891.TW","2002.TW","1301.TW",
                "0050.TW","0056.TW","00878.TW","00919.TW","2603.TW","2609.TW",
                "3711.TW","2379.TW","3008.TW","2395.TW","4938.TW","2376.TW"
            ]

        results = []
        # 解析條件關鍵字
        want_high_div   = any(k in criteria for k in ["殖利率", "配息", "dividend"])
        want_low_pe     = any(k in criteria for k in ["本益比", "PE", "pe"])
        want_high_roe   = any(k in criteria for k in ["ROE", "roe", "股東權益"])
        want_top_gain   = any(k in criteria for k in ["漲最多", "漲幅", "漲停", "上漲"])
        want_large_cap  = any(k in criteria for k in ["市值", "大型", "large"])

        for sym in candidates[:30]:  # 限制查詢數量避免超時
            try:
                info = yf.Ticker(sym).info
                name = info.get("shortName", sym)
                price = info.get("currentPrice") or info.get("regularMarketPrice", 0)
                pe    = info.get("trailingPE", 0) or 0
                div_y = (info.get("dividendYield", 0) or 0) * 100
                roe   = (info.get("returnOnEquity", 0) or 0) * 100
                cap   = info.get("marketCap", 0) or 0
                chg   = info.get("regularMarketChangePercent", 0) or 0

                score = 0
                if want_high_div and div_y > 3:   score += div_y
                if want_low_pe   and 0 < pe < 20: score += (20 - pe)
                if want_high_roe and roe > 15:     score += roe / 10
                if want_top_gain:                  score += chg
                if want_large_cap and cap > 1e11:  score += 1

                if score > 0:
                    results.append((score, sym, name, price, pe, div_y, roe, chg, cap))
            except Exception:
                pass

        results.sort(reverse=True)
        lines = [f"🔎 選股結果（{criteria}）\n"]
        for i, (_, sym, name, price, pe, div_y, roe, chg, cap) in enumerate(results[:10], 1):
            cap_str = f"{cap/1e12:.1f}兆" if cap >= 1e12 else f"{cap/1e8:.0f}億"
            line = f"{i}. {name}（{sym}）  {chg:+.1f}%"
            if div_y: line += f"  殖利率{div_y:.1f}%"
            if pe:    line += f"  PE{pe:.0f}"
            if roe:   line += f"  ROE{roe:.0f}%"
            lines.append(line)
        if not results:
            lines.append("找不到符合條件的股票，試著調整條件")
        return "\n".join(lines)
    except Exception as e:
        return f"選股篩選失敗：{e}"


def fetch_summary_writer(topic: str, max_points: int = 7, lang: str = "zh-tw") -> str:
    try:
        from ddgs import DDGS
        region = "tw-tzh" if lang == "zh-tw" else "us-en"
        all_hits = []
        with DDGS() as ddgs:
            for q in [topic, f"{topic} 重點 整理", f"{topic} 分析 摘要"]:
                hits = list(ddgs.text(q, region=region, max_results=4))
                all_hits.extend(hits)
        # 去重並提取文字
        seen, texts = set(), []
        for h in all_hits:
            title = h.get("title", "")
            body = h.get("body", "")
            if title not in seen and body:
                seen.add(title)
                texts.append(f"{title}：{body}")
        combined = "\n".join(texts[:10])
        # 從文字中提取句子當重點
        import re
        sentences = re.split(r'[。！？\n]', combined)
        sentences = [s.strip() for s in sentences if len(s.strip()) > 15][:max_points * 3]
        # 去重相似句
        points, seen_short = [], set()
        for s in sentences:
            key = s[:10]
            if key not in seen_short:
                seen_short.add(key)
                points.append(s[:120])
            if len(points) >= max_points:
                break
        lines = [f"📝 摘要：{topic}\n", f"── 核心重點（{len(points)} 項）──"]
        for i, p in enumerate(points, 1):
            lines.append(f"{i}. {p}")
        lines.append(f"\n共整合 {len(seen)} 篇資料來源")
        return "\n".join(lines)
    except Exception as e:
        return f"摘要失敗：{e}"


def fetch_timeline_events(topic: str, lang: str = "zh-tw") -> str:
    try:
        from ddgs import DDGS
        region = "tw-tzh" if lang == "zh-tw" else "us-en"
        queries = [
            f"{topic} 歷史 時間軸 發展",
            f"{topic} 大事記 年表",
            f"{topic} 起源 始末 過程",
        ]
        all_hits = []
        with DDGS() as ddgs:
            for q in queries:
                hits = list(ddgs.text(q, region=region, max_results=4))
                all_hits.extend(hits)
        lines = [f"📅 時間軸：{topic}\n"]
        seen = set()
        for h in all_hits:
            title = h.get("title", "")
            if title in seen:
                continue
            seen.add(title)
            body = h.get("body", "")[:300]
            lines.append(f"【{title}】\n{body}\n")
        lines.append(f"（共 {len(seen)} 筆資料，建議搭配 Wikipedia 查詢完整年表）")
        return "\n".join(lines)
    except Exception as e:
        return f"時間軸整理失敗：{e}"


def fetch_travel_research(destination: str, days: int = None, style: str = "", lang: str = "zh-tw") -> str:
    try:
        from ddgs import DDGS
        region = "tw-tzh" if lang == "zh-tw" else "us-en"
        style_tag = f" {style}" if style else ""
        days_tag = f" {days}天" if days else ""
        queries = {
            "景點": f"{destination} 必去景點 推薦",
            "美食": f"{destination} 必吃美食 餐廳",
            "交通住宿": f"{destination} 交通 住宿 費用",
            "注意事項": f"{destination} 旅遊注意 簽證 安全",
        }
        lines = [f"✈️ 旅遊研究：{destination}"]
        if days:
            lines.append(f"行程天數：{days} 天")
        if style:
            lines.append(f"旅遊風格：{style}")
        lines.append("")
        with DDGS() as ddgs:
            for sec, q in queries.items():
                hits = list(ddgs.text(q + style_tag + days_tag, region=region, max_results=2))
                if hits:
                    lines.append(f"── {sec} ──")
                    for h in hits:
                        lines.append(f"• {h.get('title','')}：{h.get('body','')[:200]}")
                    lines.append("")
        return "\n".join(lines)
    except Exception as e:
        return f"旅遊研究失敗：{e}"


def fetch_trend_forecast(topic: str, timeframe: str = "全部", lang: str = "zh-tw") -> str:
    try:
        from ddgs import DDGS
        region = "tw-tzh" if lang == "zh-tw" else "us-en"
        queries = [
            f"{topic} 未來趨勢 預測 展望",
            f"{topic} 短期 發展 2024 2025",
            f"{topic} 長期 影響 趨勢",
        ]
        all_hits = []
        with DDGS() as ddgs:
            for q in queries:
                hits = list(ddgs.text(q, region=region, max_results=3))
                all_hits.extend(hits)
        lines = [f"🔮 趨勢預測：{topic}", f"預測範圍：{timeframe}\n"]
        if timeframe in ("短期(1年內)", "全部"):
            lines.append("── 短期（1年內）──")
            for h in all_hits[:3]:
                lines.append(f"• {h.get('title','')}：{h.get('body','')[:180]}")
            lines.append("")
        if timeframe in ("中期(1-3年)", "全部"):
            lines.append("── 中期（1–3年）──")
            for h in all_hits[3:6]:
                lines.append(f"• {h.get('title','')}：{h.get('body','')[:180]}")
            lines.append("")
        if timeframe in ("長期(3年以上)", "全部"):
            lines.append("── 長期（3年以上）──")
            for h in all_hits[6:9]:
                lines.append(f"• {h.get('title','')}：{h.get('body','')[:180]}")
            lines.append("")
        lines.append(f"⚠️ 預測基於現有公開資料，實際發展受多重因素影響")
        return "\n".join(lines)
    except Exception as e:
        return f"趨勢預測失敗：{e}"


def fetch_tw_tax_calculator(dividend_income: float, other_income: float = 0,
                             tax_bracket: float = None, sell_amount: float = 0) -> str:
    try:
        # 健保補充保費（2.11%，超過2萬才扣）
        nhi_surcharge = dividend_income * 0.0211 if dividend_income >= 20000 else 0
        # 股利可抵減稅額（8.5%，上限8萬）
        tax_credit = min(dividend_income * 0.085, 80000)
        # 分離課稅（28%）
        separate_tax = dividend_income * 0.28
        # 合併申報
        if tax_bracket:
            total_income = dividend_income + other_income
            combined_tax = total_income * (tax_bracket / 100) - tax_credit
            combined_tax = max(combined_tax, 0)
        else:
            combined_tax = None
        # 證交稅（0.3%）
        securities_tax = sell_amount * 0.003

        lines = [
            f"💰 台股稅務試算\n",
            f"股利所得：{dividend_income:,.0f} 元",
        ]
        if other_income:
            lines.append(f"其他收入：{other_income:,.0f} 元")
        lines += [
            f"",
            f"── 健保補充保費（2.11%）──",
            f"補充保費：{nhi_surcharge:,.0f} 元",
            f"",
            f"── 方案一：分離課稅（28%）──",
            f"應繳稅額：{separate_tax:,.0f} 元",
            f"稅後股利：{dividend_income - separate_tax - nhi_surcharge:,.0f} 元",
        ]
        if combined_tax is not None:
            lines += [
                f"",
                f"── 方案二：合併申報（稅率{tax_bracket}%）──",
                f"可抵減稅額：{tax_credit:,.0f} 元",
                f"應繳稅額：{combined_tax:,.0f} 元",
                f"稅後股利：{dividend_income - combined_tax - nhi_surcharge:,.0f} 元",
                f"",
                f"建議：{'合併申報' if combined_tax < separate_tax else '分離課稅'} 節稅 {abs(separate_tax - combined_tax):,.0f} 元",
            ]
        if sell_amount > 0:
            lines += [
                f"",
                f"── 證券交易稅（0.3%）──",
                f"賣出金額：{sell_amount:,.0f} 元",
                f"證交稅：{securities_tax:,.0f} 元",
            ]
        return "\n".join(lines)
    except Exception as e:
        return f"稅務試算失敗：{e}"


def fetch_vision_locate(description: str, monitor: int = 1, action: str = "click", region: list = None) -> str:
    try:
        # 嘗試把瀏覽器切到前景並最大化（如果描述中提到影片、YouTube等）
        _desc_lower = description.lower()
        _browser_kw = any(k in _desc_lower for k in ["影片", "video", "youtube", "縮圖", "thumbnail", "搜尋結果", "第一", "chrome", "網頁"])
        if _browser_kw:
            try:
                import win32gui, win32con
                def _find_chrome(h, results):
                    if win32gui.IsWindowVisible(h):
                        t = win32gui.GetWindowText(h).lower()
                        if "youtube" in t or "chrome" in t or "edge" in t or "firefox" in t:
                            results.append((h, t))
                    return True
                _browser_wins = []
                win32gui.EnumWindows(_find_chrome, _browser_wins)
                _yt_wins = [h for h, t in _browser_wins if "youtube" in t]
                _target = _yt_wins[0] if _yt_wins else (_browser_wins[0][0] if _browser_wins else None)
                if _target:
                    win32gui.ShowWindow(_target, win32con.SW_RESTORE)
                    win32gui.ShowWindow(_target, win32con.SW_MAXIMIZE)
                    import ctypes
                    ctypes.windll.user32.SetForegroundWindow(_target)
                    import time; time.sleep(0.5)
            except Exception:
                pass

        # ── 策略1：UIA 快速搜尋（不需截圖和 API，最快） ──
        if not region and action != "locate_only":
            uia_x, uia_y = _uia_find_element(description)
            if uia_x is not None:
                if action != "locate_only":
                    # 用動作驗證迴圈點擊
                    _det_mon = monitor if not _browser_kw else (2 if monitor == 1 else monitor)
                    success = _action_with_verify(
                        lambda: _si_universal(uia_x, uia_y, action),
                        monitor=_det_mon
                    )
                    v = "已驗證" if success else "未驗證"
                    return f"✅ UIA找到「{description}」，已在 ({uia_x}, {uia_y}) 執行 {action}（{v}）"
                return f"✅ UIA找到「{description}」，位置 ({uia_x}, {uia_y})"

        # ── 策略2：YOLO 快速偵測（本地，0.05秒） ──
        if not region:
            try:
                img_yolo, mon_left_y, mon_top_y = _cap_monitor_logical(monitor)
                detections = _yolo_detect(img_yolo, conf=0.3)
                if detections:
                    # 用描述關鍵字匹配 YOLO 偵測到的物件
                    _desc_words = [w.lower() for w in description.split() if len(w) > 1]
                    best_det = None
                    best_score = 0
                    for label, cx, cy, w, h, conf in detections:
                        score = sum(1 for kw in _desc_words if kw in label.lower()) * 50 + conf * 30
                        if score > best_score:
                            best_det = (label, cx, cy, conf)
                            best_score = score
                    if best_det and best_score > 30:
                        _yl, _yx, _yy, _yc = best_det
                        abs_x = mon_left_y + _yx
                        abs_y = mon_top_y + _yy
                        if action == "locate_only":
                            return f"✅ YOLO找到「{_yl}」({_yc:.0%})，位置 ({abs_x}, {abs_y})"
                        success = _action_with_verify(
                            lambda: _si_universal(abs_x, abs_y, action),
                            monitor=monitor
                        )
                        v = "已驗證" if success else "未驗證"
                        return f"✅ YOLO找到「{_yl}」({_yc:.0%})，已在 ({abs_x}, {abs_y}) 執行 {action}（{v}）"
            except Exception:
                pass

        # ── 策略3：智慧等待 + 視覺辨識（截圖 + Sonnet） ──
        # 等螢幕穩定再截圖
        _wait_screen_stable(monitor, threshold=0.5, timeout=3.0, interval=0.3)

        img, mon_left, mon_top = _cap_monitor_logical(monitor)

        if region:
            from PIL import Image as _PI
            img = img.crop((region[0], region[1], region[0]+region[2], region[1]+region[3]))
            mon_left += region[0]; mon_top += region[1]

        rx, ry = _vision_find(img, description)
        if rx is None:
            return f"視覺辨識找不到「{description}」"

        abs_x = mon_left + rx
        abs_y = mon_top + ry

        if action == "locate_only":
            return f"✅ 找到「{description}」，位置 ({abs_x}, {abs_y})"

        # 用動作驗證迴圈點擊
        success = _action_with_verify(
            lambda: _si_universal(abs_x, abs_y, action),
            monitor=monitor
        )
        v = "已驗證" if success else "未驗證"
        return f"✅ 視覺找到「{description}」，已在 ({abs_x}, {abs_y}) 執行 {action}（{v}）"
    except Exception as e:
        return f"視覺定位失敗：{e}"


def fetch_wait_and_click(target_text: str, timeout: int = 15, monitor: int = 1, action_after: str = "click") -> str:
    try:
        import time
        start = time.time()
        interval = 1.0
        while time.time() - start < timeout:
            result = fetch_ocr_click(target_text, monitor, action_after if action_after != "none" else "click")
            if "✅" in result:
                return f"✅ 等待 {time.time()-start:.1f}s 後找到並點擊：{result}"
            if action_after == "none" and "找到" in result:
                return f"✅ 等待 {time.time()-start:.1f}s 後出現：{target_text}"
            time.sleep(interval)
        return f"⏰ 等待 {timeout}s 仍未出現「{target_text}」，超時"
    except Exception as e:
        return f"等待點擊失敗：{e}"


def fetch_warrant(underlying: str) -> str:
    try:
        # TWSE 權證查詢
        url = f"https://www.twse.com.tw/rwd/zh/warrant/MNO01?response=json&stockNo={underlying}"
        resp = requests.get(url, timeout=10, headers={"User-Agent": "Mozilla/5.0"})
        data = resp.json()

        if data.get("stat") != "OK" or not data.get("data"):
            # 改用搜尋提示
            try:
                from ddgs import DDGS
                with DDGS() as ddgs:
                    r = list(ddgs.text(f"台股 {underlying} 權證 認購 認售", region="zh-tw", max_results=3))
                lines = [f"🎫 {underlying} 相關權證資訊\n"]
                for item in r:
                    lines.append(f"• {item.get('title','')}\n  {item.get('body','')[:100]}")
                return "\n\n".join(lines)
            except Exception:
                return f"查無 {underlying} 的權證資料，請至台灣證券交易所查詢"

        rows = data.get("data", [])[:10]
        lines = [f"🎫 {underlying} 相關權證（前10筆）\n"]
        for row in rows:
            w_code = row[0]
            w_name = row[1]
            w_type = "認購" if "購" in w_name else "認售"
            strike = row[3] if len(row) > 3 else "-"
            exp = row[4] if len(row) > 4 else "-"
            lines.append(f"{w_code} {w_name}（{w_type}）　履約{strike}　到期{exp}")
        return "\n".join(lines)
    except Exception as e:
        return f"權證查詢失敗：{e}"


def fetch_weather(city: str) -> str:
    try:
        res = requests.get(
            f"https://wttr.in/{city}?format=j1",
            headers={"User-Agent": "curl/7.68.0"},
            timeout=10
        )
        data = res.json()
        current = data["current_condition"][0]
        area = data["nearest_area"][0]
        location = area["areaName"][0]["value"] + ", " + area["country"][0]["value"]
        desc = current["lang_zh-tw"][0]["value"] if "lang_zh-tw" in current else current["weatherDesc"][0]["value"]
        temp = current["temp_C"]
        feels = current["FeelsLikeC"]
        humidity = current["humidity"]
        wind = current["windspeedKmph"]
        wind_dir = current["winddir16Point"]
        return (
            f"📍 {location}\n"
            f"🌤 {desc}\n"
            f"🌡 氣溫：{temp}°C（體感 {feels}°C）\n"
            f"💧 濕度：{humidity}%\n"
            f"💨 風速：{wind} km/h（{wind_dir}）"
        )
    except Exception:
        return f"查不到「{city}」的天氣資訊。"


# ── Tavily AI 搜尋（取代 ddg_search 作為主力搜尋）──────────────────
def tavily_search(query: str, max_results: int = 5, search_depth: str = "advanced") -> str:
    """Tavily AI 搜尋 — 直接回傳乾淨的內容，不只是連結"""
    try:
        _key = os.environ.get("TAVILY_API_KEY", "")
        if not _key:
            return execute_ddg_search(query, "zh-tw", max_results)
        from tavily import TavilyClient
        client = TavilyClient(api_key=_key)
        response = client.search(query=query, search_depth=search_depth, max_results=max_results, include_answer=True)
        lines = []
        answer = response.get("answer", "")
        if answer:
            lines.append(f"📝 AI 整合答案：\n{answer}\n")
        results = response.get("results", [])
        for r in results[:max_results]:
            title = r.get("title", "")
            content = r.get("content", "")[:200]
            url = r.get("url", "")
            lines.append(f"🔍 {title}\n   {content}\n   {url}")
        return "\n\n".join(lines) if lines else "無搜尋結果"
    except Exception as e:
        logging.warning(f"Tavily failed, fallback to DDG: {e}")
        return execute_ddg_search(query, "zh-tw", max_results)


# ── ESPN 體育即時比分 ──────────────────────────────────────────────
def fetch_sports_scores(sport: str = "nba", league: str = "", date: str = "") -> str:
    """ESPN 即時比分 — 支援 NBA/NFL/MLB/NHL/足球"""
    try:
        sport_map = {
            "nba": ("basketball", "nba"), "wnba": ("basketball", "wnba"),
            "nfl": ("football", "nfl"), "mlb": ("baseball", "mlb"),
            "nhl": ("hockey", "nhl"), "mls": ("soccer", "usa.1"),
            "epl": ("soccer", "eng.1"), "laliga": ("soccer", "esp.1"),
            "bundesliga": ("soccer", "ger.1"), "seriea": ("soccer", "ita.1"),
            "ligue1": ("soccer", "fra.1"), "champions": ("soccer", "uefa.champions"),
            "worldcup": ("soccer", "fifa.world"),
        }
        key = sport.lower().replace(" ", "")
        if key not in sport_map:
            return f"不支援「{sport}」，支援的有：{', '.join(sport_map.keys())}"
        sport_cat, league_id = sport_map[key]
        url = f"https://site.api.espn.com/apis/site/v2/sports/{sport_cat}/{league_id}/scoreboard"
        if date:
            url += f"?dates={date.replace('-', '')}"
        resp = requests.get(url, timeout=10, headers={"User-Agent": "Mozilla/5.0"})
        data = resp.json()
        events = data.get("events", [])
        if not events:
            return f"🏟️ {sport.upper()} 今天沒有比賽"
        sport_emoji = {"nba": "🏀", "nfl": "🏈", "mlb": "⚾", "nhl": "🏒"}.get(key, "⚽")
        lines = [f"{sport_emoji} {sport.upper()} 比分\n"]
        for e in events:
            status_desc = e.get("status", {}).get("type", {}).get("description", "")
            comps = e.get("competitions", [{}])[0]
            teams = comps.get("competitors", [])
            if len(teams) >= 2:
                away, home = teams[0], teams[1]
                away_name = away["team"]["displayName"]
                away_score = away.get("score", "?")
                home_name = home["team"]["displayName"]
                home_score = home.get("score", "?")
                series = comps.get("series", {})
                series_info = f"  📊 {series.get('summary', '')}" if series.get("summary") else ""
                if status_desc == "Final":
                    lines.append(f"  {away_name} {away_score} - {home_score} {home_name}  🏁{series_info}")
                elif status_desc == "Scheduled":
                    game_time = e.get("status", {}).get("type", {}).get("detail", "")
                    lines.append(f"  {away_name} vs {home_name}  ⏰ {game_time}{series_info}")
                else:
                    period = e.get("status", {}).get("period", "")
                    clock = e.get("status", {}).get("displayClock", "")
                    lines.append(f"  {away_name} {away_score} - {home_score} {home_name}  🔴 LIVE Q{period} {clock}{series_info}")
        return "\n".join(lines)
    except Exception as e:
        return f"體育比分查詢失敗：{e}"


# ── NewsAPI 全球新聞 ──────────────────────────────────────────────
def fetch_global_news(query: str = "", category: str = "", country: str = "", count: int = 5) -> str:
    """NewsAPI 全球新聞搜尋 — 80000+ 來源"""
    try:
        _key = os.environ.get("NEWSAPI_KEY", "")
        if not _key:
            return search_news(query or "top news", "zh-TW", count)
        count = min(max(count, 1), 10)
        if query:
            url = f"https://newsapi.org/v2/everything?q={requests.utils.quote(query)}&sortBy=publishedAt&pageSize={count}&apiKey={_key}"
        elif category:
            cat_map = {"科技": "technology", "商業": "business", "體育": "sports",
                       "娛樂": "entertainment", "健康": "health", "科學": "science"}
            cat = cat_map.get(category, category.lower())
            ctry = country or "us"
            url = f"https://newsapi.org/v2/top-headlines?category={cat}&country={ctry}&pageSize={count}&apiKey={_key}"
        else:
            ctry = country or "us"
            url = f"https://newsapi.org/v2/top-headlines?country={ctry}&pageSize={count}&apiKey={_key}"
        resp = requests.get(url, timeout=10)
        data = resp.json()
        if data.get("status") != "ok":
            return f"新聞查詢失敗：{data.get('message', 'unknown error')}"
        articles = data.get("articles", [])
        if not articles:
            return "沒有找到相關新聞"
        lines = [f"📰 全球新聞" + (f"（{query}）" if query else "") + "\n"]
        for a in articles[:count]:
            title = a.get("title", "")
            source = a.get("source", {}).get("name", "")
            published = a.get("publishedAt", "")[:16].replace("T", " ")
            desc = a.get("description", "") or ""
            lines.append(f"• {title}\n  {source} | {published}\n  {desc[:100]}")
        return "\n\n".join(lines)
    except Exception as e:
        return f"新聞查詢失敗：{e}"


# ── 百度搜尋（中國資訊專用）──────────────────────────────────────
def baidu_search(query: str, max_results: int = 5) -> str:
    """搜尋中國大陸的資訊 — 用 DuckDuckGo cn-zh 區域 + 中國新聞源"""
    try:
        from ddgs import DDGS
        lines = []
        with DDGS() as ddgs:
            results = list(ddgs.text(query, region="cn-zh", max_results=max_results))
            for r in results:
                title = r.get("title", "")
                body = r.get("body", "")[:150]
                href = r.get("href", "")
                lines.append(f"🔍 {title}\n   {body}\n   {href}")
        try:
            import feedparser
            feed_url = f"https://news.google.com/rss/search?q={requests.utils.quote(query)}&hl=zh-CN&gl=CN&ceid=CN:zh-Hans"
            feed = feedparser.parse(feed_url)
            for entry in feed.entries[:3]:
                title = entry.get("title", "").split(" - ")[0]
                pub = entry.get("published", "")[:16]
                lines.append(f"📰 {title}（{pub}）")
        except Exception:
            pass
        if not lines:
            return tavily_search(query, max_results)
        return "\n\n".join(lines)
    except Exception as e:
        return f"百度搜尋失敗：{e}"


def fetch_window_manager(action: str = "list", window_name: str = "") -> str:
    """視窗管理：列出所有視窗 / 切換焦點 / 最大化 / 最小化 / 關閉"""
    try:
        import win32gui, win32con, re
        results = []

        if action == "list":
            wins = []
            def _enum(h, _):
                if win32gui.IsWindowVisible(h):
                    t = win32gui.GetWindowText(h).strip()
                    if t: wins.append(f"[{h}] {t}")
                return True
            win32gui.EnumWindows(_enum, None)
            return "開啟中的視窗：\n" + "\n".join(wins[:30])

        # 找目標視窗
        matches = []
        def _find(h, _):
            if win32gui.IsWindowVisible(h):
                t = win32gui.GetWindowText(h)
                if window_name.lower() in t.lower():
                    matches.append(h)
            return True
        win32gui.EnumWindows(_find, None)

        if not matches:
            return f"找不到包含「{window_name}」的視窗"

        hw = matches[0]
        title = win32gui.GetWindowText(hw)

        if action == "focus":
            win32gui.ShowWindow(hw, win32con.SW_RESTORE)
            win32gui.SetForegroundWindow(hw)
            return f"✅ 已切換到視窗：{title}"
        elif action == "maximize":
            win32gui.ShowWindow(hw, win32con.SW_MAXIMIZE)
            return f"✅ 已最大化：{title}"
        elif action == "minimize":
            win32gui.ShowWindow(hw, win32con.SW_MINIMIZE)
            return f"✅ 已最小化：{title}"
        elif action == "close":
            win32gui.PostMessage(hw, win32con.WM_CLOSE, 0, 0)
            return f"✅ 已關閉：{title}"
        else:
            return f"未知動作：{action}，可用：list/focus/maximize/minimize/close"
    except Exception as e:
        return f"視窗管理失敗：{e}"


def generate_candlestick(symbol: str, period: str = "3mo") -> tuple[bytes | None, str]:
    """回傳 (PNG bytes, 型態分析文字)"""
    try:
        import yfinance as yf
        import mplfinance as mpf
        import tempfile, ta as ta_lib
        ticker = yf.Ticker(symbol)
        hist = ticker.history(period=period)
        if hist.empty:
            return None, f"找不到「{symbol}」資料"

        name = ticker.info.get("shortName") or symbol
        close = hist["Close"]
        ma20 = close.rolling(20).mean()
        ma60 = close.rolling(60).mean()

        add_plots = [
            mpf.make_addplot(ma20, color="orange", width=1.2, label="MA20"),
            mpf.make_addplot(ma60, color="purple", width=1.2, label="MA60"),
        ]

        with tempfile.NamedTemporaryFile(suffix=".png", delete=False) as f:
            out_path = f.name

        mpf.plot(
            hist, type="candle", style="charles",
            title=f"{name} ({symbol})",
            ylabel="Price", volume=True,
            addplot=add_plots,
            savefig=dict(fname=out_path, dpi=150, bbox_inches="tight"),
            figsize=(12, 7)
        )

        # 型態辨識
        patterns = []
        c = close.values
        n = len(c)
        if n >= 20:
            recent = c[-20:]
            peak_idx = recent.argmax()
            trough_idx = recent.argmin()
            # 簡單型態判斷
            if c[-1] > ma20.iloc[-1] > ma60.iloc[-1]:
                patterns.append("多頭排列（MA20>MA60，趨勢向上）📈")
            elif c[-1] < ma20.iloc[-1] < ma60.iloc[-1]:
                patterns.append("空頭排列（MA20<MA60，趨勢向下）📉")
            if peak_idx < 5 and c[-1] < recent[peak_idx] * 0.95:
                patterns.append("近期高點已過，回落中")
            if trough_idx < 5 and c[-1] > recent[trough_idx] * 1.05:
                patterns.append("近期低點反彈，留意支撐")
            # 突破判斷
            resistance = max(c[-20:-5])
            support = min(c[-20:-5])
            if c[-1] > resistance:
                patterns.append(f"突破近期壓力 {resistance:.2f} ⚡")
            elif c[-1] < support:
                patterns.append(f"跌破近期支撐 {support:.2f} ⚠️")

        pattern_str = "\n".join(patterns) if patterns else "無明顯型態訊號"

        with open(out_path, "rb") as f:
            img_bytes = f.read()
        Path(out_path).unlink(missing_ok=True)
        return img_bytes, pattern_str
    except Exception as e:
        return None, f"K線圖生成失敗：{e}"


def generate_voice_ogg(text: str, voice: str = "zh-CN-YunxiNeural") -> bytes:
    """生成語音並回傳 OGG OPUS bytes（Telegram voice message 格式）
    優先使用 XTTS v2（更自然人聲），失敗自動 fallback 到 edge_tts
    """
    text = clean_for_tts(text)
    import tempfile, subprocess as sp
    import imageio_ffmpeg

    ffmpeg_exe = imageio_ffmpeg.get_ffmpeg_exe()
    tmp_ogg = tempfile.NamedTemporaryFile(suffix=".ogg", delete=False)
    tmp_ogg.close()

    # ── 嘗試 XTTS v2 ──────────────────────────────────────
    wav_bytes = None
    if _ensure_xtts_server():
        wav_bytes = _xtts_generate_wav(text)

    if wav_bytes:
        # XTTS WAV → OGG OPUS（加低音 EQ）
        tmp_wav = tempfile.NamedTemporaryFile(suffix=".wav", delete=False)
        tmp_wav.write(wav_bytes)
        tmp_wav.close()
        sp.run([
            ffmpeg_exe, "-y", "-i", tmp_wav.name,
            "-af", "equalizer=f=80:width_type=o:width=2:g=6,equalizer=f=150:width_type=o:width=2:g=4",
            "-c:a", "libopus", "-b:a", "96k",
            tmp_ogg.name
        ], capture_output=True)
        Path(tmp_wav.name).unlink(missing_ok=True)
    else:
        # ── Fallback：edge_tts ──────────────────────────────
        import edge_tts, asyncio
        tmp_mp3 = tempfile.NamedTemporaryFile(suffix=".mp3", delete=False)
        tmp_mp3.close()

        async def _gen():
            comm = edge_tts.Communicate(text, voice, rate="-10%", pitch="-15Hz")
            await comm.save(tmp_mp3.name)
        asyncio.run(_gen())

        # MP3 → OGG OPUS（加低音 EQ）
        sp.run([
            ffmpeg_exe, "-y", "-i", tmp_mp3.name,
            "-af", "equalizer=f=80:width_type=o:width=2:g=6,equalizer=f=150:width_type=o:width=2:g=4",
            "-c:a", "libopus", "-b:a", "96k",
            tmp_ogg.name
        ], capture_output=True)
        Path(tmp_mp3.name).unlink(missing_ok=True)

    data = Path(tmp_ogg.name).read_bytes()
    Path(tmp_ogg.name).unlink(missing_ok=True)
    return data


def init_db():
    with _db_lock:
        conn = _get_db()
        conn.execute("""
            CREATE TABLE IF NOT EXISTS chat_history (
                id INTEGER PRIMARY KEY AUTOINCREMENT,
                chat_id INTEGER NOT NULL,
                role TEXT NOT NULL,
                content TEXT NOT NULL,
                created_at TIMESTAMP DEFAULT CURRENT_TIMESTAMP
            )
        """)
        conn.execute("""
            CREATE TABLE IF NOT EXISTS long_term_memory (
                id INTEGER PRIMARY KEY AUTOINCREMENT,
                chat_id INTEGER NOT NULL,
                content TEXT NOT NULL,
                created_at TIMESTAMP DEFAULT CURRENT_TIMESTAMP
            )
        """)
        conn.commit()


def log_message(direction: str, sender: str, chat_id: int, text: str):
    """寫入訊息日誌供終端機同步使用"""
    ts = dt.dt.datetime.now().strftime("%H:%M:%S")
    line = f"[{ts}] {direction} [{chat_id}] {sender}: {text}\n"
    with open(MSG_LOG, "a", encoding="utf-8") as f:
        f.write(line)


def multi_perspective(topic: str, lang: str = "zh-tw") -> str:
    try:
        from ddgs import DDGS
        results = {}
        queries = {
            "支持/正面觀點": f"{topic} 優點 支持 正面",
            "反對/批評觀點": f"{topic} 缺點 反對 批評 問題",
            "中立/分析觀點": f"{topic} 分析 評估 影響 研究",
        }
        if lang == "en":
            queries = {
                "Pro / Positive": f"{topic} benefits support positive",
                "Con / Critical": f"{topic} criticism problems negative against",
                "Neutral / Analysis": f"{topic} analysis impact research objective",
            }
        with DDGS() as ddgs:
            for label, q in queries.items():
                items = list(ddgs.text(q, region=lang, max_results=3))
                lines = [f"── {label} ──"]
                for r in items:
                    lines.append(f"• {r['title']}")
                    lines.append(f"  {r['body'][:150]}")
                results[label] = "\n".join(lines)
        output = [f"🔍 多角度分析：{topic}\n"]
        output.extend(results.values())
        return "\n\n".join(output)
    except Exception as e:
        return f"多角度分析失敗：{e}"


def ptt_search(keyword: str, board: str = "Gossiping", count: int = 5) -> str:
    try:
        from bs4 import BeautifulSoup
        import ssl, urllib3
        count = min(count, 10)
        # PTT 需要 session + 放寬 SSL + Cookie
        session = requests.Session()
        session.verify = False
        urllib3.disable_warnings(urllib3.exceptions.InsecureRequestWarning)
        headers = {
            "User-Agent": "Mozilla/5.0 (Windows NT 10.0; Win64; x64) AppleWebKit/537.36 Chrome/120.0",
            "Cookie": "over18=1",
            "Referer": "https://www.ptt.cc/",
        }
        search_url = f"https://www.ptt.cc/bbs/{board}/search?q={urllib.parse.quote(keyword)}"
        resp = session.get(search_url, headers=headers, timeout=12)
        soup = BeautifulSoup(resp.text, "html.parser")
        posts = soup.select(".r-ent")[:count]
        if not posts:
            # fallback: 用 Google News 搜尋 PTT 相關文章
            import feedparser
            fallback_url = f"https://news.google.com/rss/search?q=PTT+{urllib.parse.quote(keyword)}&hl=zh-Hant&gl=TW&ceid=TW:zh-Hant"
            feed = feedparser.parse(fallback_url)
            if feed.entries:
                lines = [f"📋 PTT 搜尋結果（Google News 補充）：{keyword}\n"]
                for i, e in enumerate(feed.entries[:count], 1):
                    lines.append(f"{i}. {e.get('title', '')}")
                return "\n".join(lines)
            return f"PTT {board} 版找不到「{keyword}」相關文章"
        lines = [f"📋 PTT/{board} 搜尋：{keyword}\n"]
        for post in posts:
            title_el = post.select_one(".title a")
            meta_el = post.select_one(".meta .author")
            date_el = post.select_one(".meta .date")
            nrec_el = post.select_one(".nrec span")
            if not title_el:
                continue
            title = title_el.get_text(strip=True)
            author = meta_el.get_text(strip=True) if meta_el else ""
            date = date_el.get_text(strip=True) if date_el else ""
            nrec = nrec_el.get_text(strip=True) if nrec_el else "0"
            post_url = "https://www.ptt.cc" + title_el["href"]
            lines.append(f"🗂 {title}")
            lines.append(f"   推文：{nrec}　作者：{author}　{date}")
            try:
                p_resp = session.get(post_url, headers=headers, timeout=6)
                p_soup = BeautifulSoup(p_resp.text, "html.parser")
                content_el = p_soup.select_one("#main-content")
                if content_el:
                    for tag in content_el.select(".f2, .push, #article-polling"):
                        tag.decompose()
                    content = content_el.get_text(separator=" ", strip=True)[:200]
                    lines.append(f"   摘要：{content}")
            except Exception:
                pass
        return "\n".join(lines)
    except Exception as e:
        return f"PTT 搜尋失敗：{e}"


def read_webpage(url: str, max_chars: int = 3000) -> str:
    """讀取網頁內容，先用 requests，內容太少或失敗時自動 fallback 到 Playwright"""
    def _parse_html(html, url):
        from bs4 import BeautifulSoup
        soup = BeautifulSoup(html, "html.parser")
        for tag in soup(["script", "style", "nav", "footer", "header", "aside", "ads"]):
            tag.decompose()
        content_el = soup.find("article") or soup.find("main") or soup.find("body")
        text = content_el.get_text(separator="\n", strip=True) if content_el else soup.get_text(separator="\n", strip=True)
        lines = [l.strip() for l in text.splitlines() if l.strip()]
        clean = "\n".join(lines)
        title = soup.title.string.strip() if soup.title else url
        return title, clean

    # 策略1：requests（快）
    try:
        headers = {"User-Agent": "Mozilla/5.0 (Windows NT 10.0; Win64; x64) AppleWebKit/537.36"}
        resp = requests.get(url, headers=headers, timeout=12)
        resp.encoding = resp.apparent_encoding
        title, clean = _parse_html(resp.text, url)
        # 如果內容太少（<100字），可能是動態渲染頁面，fallback 到 Playwright
        if len(clean) >= 100:
            result = f"【{title}】\n{url}\n\n{clean[:max_chars]}"
            if len(clean) > max_chars:
                result += f"\n\n（內容已截斷，共 {len(clean)} 字）"
            return result
    except Exception:
        pass

    # 策略2：Playwright headless（能渲染 JavaScript 動態頁面）
    try:
        import subprocess, sys, json as _json_rw, textwrap
        _script = textwrap.dedent(f"""
import sys, json
from playwright.sync_api import sync_playwright
with sync_playwright() as p:
    b = p.chromium.launch(headless=True)
    page = b.new_page(user_agent="Mozilla/5.0 (Windows NT 10.0; Win64; x64) AppleWebKit/537.36")
    page.goto({_json_rw.dumps(url)}, wait_until="networkidle", timeout=20000)
    page.wait_for_timeout(2000)
    html = page.content()
    b.close()
    print(html)
""")
        proc = subprocess.run(
            [sys.executable, "-c", _script],
            capture_output=True, text=True, encoding="utf-8", errors="replace",
            timeout=30, env={**__import__("os").environ, "PYTHONIOENCODING": "utf-8"}
        )
        if proc.returncode == 0 and proc.stdout.strip():
            title, clean = _parse_html(proc.stdout, url)
            if clean:
                result = f"【{title}】（Playwright）\n{url}\n\n{clean[:max_chars]}"
                if len(clean) > max_chars:
                    result += f"\n\n（內容已截斷，共 {len(clean)} 字）"
                return result
    except Exception:
        pass

    return f"網頁讀取失敗：無法透過 requests 或 Playwright 取得 {url} 的內容"


def search_news(query: str, lang: str = "zh-TW", count: int = 6) -> str:
    try:
        import feedparser
        count = min(count, 10)
        lang_map = {"zh-TW": "zh-Hant&gl=TW&ceid=TW:zh-Hant",
                    "zh-CN": "zh-Hans&gl=CN&ceid=CN:zh-Hans",
                    "en-US": "en&gl=US&ceid=US:en"}
        param = lang_map.get(lang, lang_map["zh-TW"])
        url = f"https://news.google.com/rss/search?q={urllib.parse.quote(query)}&hl={param}"
        feed = feedparser.parse(url)
        if not feed.entries:
            return f"找不到「{query}」的新聞"
        lines = [f"📰 {query} 最新新聞\n"]
        for i, entry in enumerate(feed.entries[:count], 1):
            title = entry.get("title", "無標題")
            pub = entry.get("published", "")[:16]
            lines.append(f"{i}. {title}（{pub}）")
        return "\n".join(lines)
    except Exception as e:
        return f"新聞搜尋失敗：{e}"


def wikipedia_search(query: str, lang: str = "zh") -> str:
    try:
        search_url = f"https://{lang}.wikipedia.org/w/api.php"
        # 先搜尋
        params = {"action": "search", "list": "search", "srsearch": query,
                  "format": "json", "srlimit": 1}
        resp = requests.get(search_url, params=params, timeout=10)
        results = resp.json().get("query", {}).get("search", [])
        if not results:
            return f"Wikipedia 找不到「{query}」相關條目"
        title = results[0]["title"]
        # 取得摘要
        params2 = {"action": "query", "titles": title, "prop": "extracts",
                   "exintro": True, "explaintext": True, "format": "json"}
        resp2 = requests.get(search_url, params=params2, timeout=10)
        pages = resp2.json().get("query", {}).get("pages", {})
        page = next(iter(pages.values()))
        extract = page.get("extract", "無法取得內容")
        extract = extract[:2500]
        return f"📖 Wikipedia：{title}\n\n{extract}"
    except Exception as e:
        return f"Wikipedia 查詢失敗：{e}"


class OperationStateMachine:
    """多步驟操作的狀態機，每一步有成功/失敗判斷和重試"""
    def __init__(self, name="operation"):
        self.name = name
        self.steps = []
        self.current = 0
        self.results = []

    def add_step(self, name: str, action, verify=None, max_retries=3, on_fail="retry"):
        self.steps.append({"name": name, "action": action, "verify": verify,
                          "max_retries": max_retries, "on_fail": on_fail})

    def run(self) -> dict:
        self.current = 0
        self.results = []
        for i, step in enumerate(self.steps):
            self.current = i
            success = False
            result = None
            for attempt in range(step["max_retries"]):
                try:
                    result = step["action"]()
                    if step["verify"] is None:
                        success = True
                        break
                    elif step["verify"](result):
                        success = True
                        break
                except Exception as e:
                    result = str(e)
                import time; time.sleep(0.5)
            self.results.append({"step": step["name"], "success": success, "result": result})
            if not success:
                if step["on_fail"] == "abort":
                    return {"ok": False, "results": self.results, "failed_step": step["name"]}
                elif step["on_fail"] == "skip":
                    continue
                else:
                    return {"ok": False, "results": self.results, "failed_step": step["name"]}
        return {"ok": True, "results": self.results, "failed_step": None}


def youtube_play_flow(keyword: str, monitor: int = 2) -> dict:
    """YouTube 搜尋並播放的完整狀態機流程"""
    import webbrowser, time as _t_yt

    sm = OperationStateMachine("youtube_play")

    # Step 1: 開 YouTube 搜尋頁
    def _open():
        webbrowser.open(f"https://www.youtube.com/results?search_query={keyword.replace(' ', '+')}")
        return True
    sm.add_step("open_search", _open, max_retries=2, on_fail="abort")

    # Step 2: 等頁面載入
    def _wait_load():
        _t_yt.sleep(1)  # 給瀏覽器反應時間
        return _wait_screen_stable(monitor, threshold=0.5, timeout=10.0)
    def _verify_load(result):
        return result  # True = 穩定了
    sm.add_step("wait_load", _wait_load, _verify_load, max_retries=2, on_fail="abort")

    # Step 3: 滾過廣告
    def _scroll():
        import pyautogui, win32gui, win32con, ctypes
        wins = []
        def _fb(h, _):
            if win32gui.IsWindowVisible(h):
                t = win32gui.GetWindowText(h).lower()
                if 'youtube' in t or 'chrome' in t:
                    wins.append(h)
            return True
        win32gui.EnumWindows(_fb, None)
        if wins:
            ctypes.windll.user32.SetForegroundWindow(wins[0])
            _t_yt.sleep(0.5)
        pyautogui.scroll(-3)
        _t_yt.sleep(1)
        return True
    sm.add_step("scroll_past_ads", _scroll, max_retries=1, on_fail="skip")

    # Step 4: vision_locate 點擊影片
    def _click_video():
        return fetch_vision_locate(
            f"YouTube搜尋結果中{keyword}的官方MV或歌曲影片縮圖（有OFFICIAL標誌的優先，跳過贊助商廣告）",
            monitor, "click"
        )
    def _verify_click(result):
        return "找到" in str(result) and "找不到" not in str(result)
    sm.add_step("click_video", _click_video, _verify_click, max_retries=3, on_fail="abort")

    # Step 5: 等影片開始播放
    def _wait_play():
        _t_yt.sleep(1)
        return _wait_screen_stable(monitor, threshold=0.3, timeout=5.0)
    sm.add_step("wait_play", _wait_play, max_retries=1, on_fail="skip")

    return sm.run()


def youtube_summary(url: str) -> str:
    try:
        from youtube_transcript_api import YouTubeTranscriptApi
        import re
        # 從 URL 提取 video ID
        patterns = [
            r"(?:v=|youtu\.be/|/embed/|/v/)([a-zA-Z0-9_-]{11})",
        ]
        video_id = url if re.match(r"^[a-zA-Z0-9_-]{11}$", url) else None
        if not video_id:
            for pat in patterns:
                m = re.search(pat, url)
                if m:
                    video_id = m.group(1)
                    break
        if not video_id:
            return f"無法從網址提取 YouTube 影片 ID：{url}"
        # 嘗試取得字幕（優先繁中→簡中→英文）
        transcript = None
        for lang in [["zh-TW", "zh-Hant"], ["zh", "zh-Hans"], ["en"]]:
            try:
                transcript = YouTubeTranscriptApi.get_transcript(video_id, languages=lang)
                break
            except Exception:
                continue
        if transcript is None:
            try:
                transcript = YouTubeTranscriptApi.get_transcript(video_id)
            except Exception as e:
                return f"無法取得字幕（影片可能無字幕或受限）：{e}"
        # 合併字幕文字
        full_text = " ".join(t["text"] for t in transcript)
        # 截取前 3000 字
        summary = full_text[:3000]
        if len(full_text) > 3000:
            summary += f"\n\n（字幕共 {len(full_text)} 字，已截取前段）"
        return f"🎬 YouTube 字幕摘要\n影片ID：{video_id}\n\n{summary}"
    except Exception as e:
        return f"YouTube 字幕擷取失敗：{e}"

# ── 缺口1：觸發驅動 ─────────────────────────────────

def email_trigger(action, host="", user="", password="", filter_from="",
                  filter_subject="", duration=300, to="", subject="", body=""):
    print(execute_email_trigger(action, host, user, password, filter_from,
                                filter_subject, duration, to, subject, body))

def file_trigger(folder, event, action, pattern="", target="", duration=60):
    print(execute_file_trigger(folder, event, action, pattern, target, duration))

def webhook_server(action, port=8765, secret=""):
    print(execute_webhook_server(action, int(port), secret))

# ── 缺口2：應用程式深度控制 ──────────────────────────

def com_auto(app, action, path="", sheet=None, cell="", value="", macro="", to="", subject=""):
    print(execute_com_auto(app, action, path, sheet, cell, value, macro, to, subject))

def dialog_auto(action, button_text="", window_title="", timeout=30):
    print(execute_dialog_auto(action, button_text, window_title, int(timeout)))

def ime_switch(action):
    print(execute_ime_switch(action))

# ── 缺口3：感知能力 ──────────────────────────────────

def wake_word(action, keyword="", duration=5, language="zh-TW"):
    print(execute_wake_word(action, keyword, float(duration), language))

def sound_detect(action, threshold=20, duration=5, output=""):
    print(execute_sound_detect(action, float(threshold), float(duration), output))

def face_recognize(action, name="", image_path="", output=""):
    print(execute_face_recognize(action, name, image_path, output))

# ── 缺口4：跨裝置控制 ────────────────────────────────

def http_server(action, port=9876, password=""):
    print(execute_http_server(action, int(port), password))

def lan_scan(action, subnet="", host="", port=80):
    print(execute_lan_scan(action, subnet, host, int(port)))

def serial_port(action, port="", baudrate=9600, data="", timeout=2):
    print(execute_serial_port(action, port, int(baudrate), data, float(timeout)))

def mqtt(action, broker, port=1883, topic="", message="", duration=10, username="", password=""):
    print(execute_mqtt(action, broker, int(port), topic, message, float(duration), username, password))

# ── 缺口5：內容理解與處理 ────────────────────────────

def doc_ai(action, path="", path2="", fields="", question="", url=""):
    print(execute_doc_ai(action, path, path2, fields, question, url))


def audio_transcribe(action, path="", duration=30, language="", output=""):
    print(execute_audio_transcribe(action, path, float(duration), language, output))

# ══════════════════════════════════════════════════════
# 奧創升級技能集
# ══════════════════════════════════════════════════════

def osint_search(action, query="", target="", limit=10):
    print(execute_osint_search(action, query, target, int(limit)))

def news_monitor(action, keywords="", interval=300, duration=3600):
    print(execute_news_monitor(action, keywords, float(interval), float(duration)))

def threat_intel(action, target="", api_key=""):
    print(execute_threat_intel(action, target, api_key))

def auto_skill(action, goal="", skill_name="", code=""):
    print(execute_auto_skill(action, goal, skill_name, code))

def smart_home(action, device="", value="", host="", token=""):
    print(execute_smart_home(action, device, value, host, token))

def goal_manager(action, goal="", goal_id="", steps="", priority="normal"):
    print(execute_goal_manager(action, goal, goal_id, steps, priority))

def auto_trade(action, symbol="", amount=0.0, price=0.0, order_type="market"):
    print(execute_auto_trade(action, symbol, float(amount), float(price), order_type))

def knowledge_base(action, content="", query="", tag="", kb_id=""):
    print(execute_knowledge_base(action, content, query, tag, kb_id))

def emotion_detect(action, text="", image_path=""):
    print(execute_emotion_detect(action, text, image_path))

def voice_id(action, name="", audio_path="", duration=5):
    print(execute_voice_id(action, name, audio_path, float(duration)))

def pentest(action, target="", port_range="1-1000", timeout=2):
    print(execute_pentest(action, target, port_range, float(timeout)))

def proactive_alert(action, name="", condition="", threshold="", target="", interval=60):
    print(execute_proactive_alert(action, name, condition, threshold, target, float(interval)))

def multi_deploy(action, remote_host="", remote_user="", remote_pass="", remote_path="/tmp/niu_bot"):
    print(execute_multi_deploy(action, remote_host, remote_user, remote_pass, remote_path))

def self_benchmark(action):
    print(execute_self_benchmark(action))


# ── 新增投資技能 ─────────────────────────────────────

def get_institutional(symbol="", date=""):
    """台股三大法人買賣超。symbol=股票代號（空=整體市場），date=YYYYMMDD（空=今天）"""
    print(fetch_institutional(symbol, date))

def get_sector(market="us"):
    """產業類股表現。market=us（美股）或 tw（台股）"""
    print(fetch_sector(market))

def get_commodity(items="all"):
    """大宗商品報價。items=gold,oil,silver,copper,natgas,wheat,corn 或 all"""
    item_list = [x.strip() for x in items.split(",")] if items != "all" else ["all"]
    print(fetch_commodity(item_list))

def get_bond_yield():
    """美國公債殖利率（2Y/5Y/10Y/30Y）及利差分析"""
    print(fetch_bond_yield())

def get_dividend_calendar(symbol):
    """除權息資訊。symbol=股票代號（如 0056.TW、AAPL）"""
    print(fetch_dividend_calendar(symbol))

def stock_screener(criteria, market="us"):
    """選股篩選。criteria=篩選條件（如「殖利率>5%」），market=us/tw"""
    print(fetch_stock_screener(criteria, market))

def get_margin_trading(symbol, date=""):
    """台股融資融券餘額。symbol=台股代號，date=YYYYMMDD（空=今天）"""
    print(fetch_margin_trading(symbol, date))

def get_options(symbol, expiry=""):
    """選擇權鏈。symbol=股票代號（如 AAPL），expiry=到期日（空=最近一個）"""
    print(fetch_options(symbol, expiry))

def get_futures(items="all"):
    """期貨報價。items=sp500,nasdaq,dow,gold,oil,taiex 或 all"""
    item_list = [x.strip() for x in items.split(",")] if items != "all" else ["all"]
    print(fetch_futures(item_list))

def get_ipo(count=10):
    """近期 IPO 行事曆。count=顯示筆數（預設10）"""
    print(fetch_ipo(int(count)))

def backtest(symbol, strategy="ma_cross", period="2y"):
    """回測投資策略。strategy=ma_cross/buy_hold/dca，period=1y/2y/3y/5y"""
    print(fetch_backtest(symbol, strategy, period))

def get_ashare(code, period="1mo"):
    """A股/港股查詢。code=6位A股代號或4位港股代號，period=1mo/3mo/6mo/1y"""
    print(fetch_ashare(code, period))

def get_cn_news(source="all", count=5):
    """中國大陸新聞。source=xinhua/people/36kr/caixin/all，count=顯示則數"""
    print(fetch_cn_news(source, int(count)))

def china_search(query, category="其他", count=6):
    """中國大陸全方位搜尋（旅遊/美食/文化/戲劇/演員/工作等）。"""
    print(fetch_china_search(query, category, int(count)))

def get_global_market():
    """全球主要股市指數概覽"""
    print(fetch_global_market())

def get_economic_calendar(count=10):
    """重要經濟數據行事曆（CPI/非農/GDP/Fed）"""
    print(fetch_economic_calendar(int(count)))

def get_earnings_calendar(days=7):
    """未來N天財報日曆"""
    print(fetch_earnings_calendar(int(days)))

def get_analyst_ratings(symbol):
    """分析師評級升降評紀錄。symbol=股票代號"""
    print(fetch_analyst_ratings(symbol))

def get_short_interest(symbol):
    """空頭比率/借券賣出資料。symbol=股票代號"""
    print(fetch_short_interest(symbol))

def get_correlation(symbols, period="1y"):
    """多股相關性矩陣。symbols=逗號分隔代號，period=3mo/6mo/1y/2y"""
    sym_list = [s.strip() for s in symbols.split(",")] if isinstance(symbols, str) else symbols
    print(fetch_correlation(sym_list, period))

def get_risk_metrics(symbol, period="1y"):
    """風險指標：Beta/夏普/波動率/VaR。period=1y/2y/3y"""
    print(fetch_risk_metrics(symbol, period))

def get_money_flow(symbol):
    """個股資金流向分析。symbol=股票代號"""
    print(fetch_money_flow(symbol))

def get_concept_stocks(theme):
    """台股概念股查詢。theme=AI/電動車/軍工/低軌衛星/半導體等"""
    print(fetch_concept_stocks(theme))

def get_crypto_depth(coin="bitcoin"):
    """加密幣深度：鏈上數據/資金費率/DeFi。coin=bitcoin/ethereum/solana等"""
    print(fetch_crypto_depth(coin))

def drip_calculator(symbol, shares, years=10, monthly_invest=0):
    """DRIP股息再投資試算。shares=初始股數，years=持有年數，monthly_invest=每月追加"""
    print(fetch_drip_calculator(symbol, float(shares), int(years), float(monthly_invest)))

def get_forex_chart(pair, period="3mo"):
    """外匯技術分析。pair=USDTWD=X等，period=1mo/3mo/6mo/1y"""
    print(fetch_forex_chart(pair, period))

def get_warrant(underlying):
    """台股認購/認售權證。underlying=標的股代號（如2330）"""
    print(fetch_warrant(underlying))

def get_portfolio_risk(symbols_weights, period="1y"):
    """投資組合風險分析。symbols_weights=代號:權重,代號:權重（如AAPL:0.5,MSFT:0.5）"""
    if isinstance(symbols_weights, str):
        holdings = []
        for item in symbols_weights.split(","):
            parts = item.strip().split(":")
            if len(parts) == 2:
                holdings.append({"symbol": parts[0].strip(), "weight": float(parts[1].strip())})
            else:
                holdings.append({"symbol": parts[0].strip(), "weight": 1.0})
    else:
        holdings = symbols_weights
    print(fetch_portfolio_risk(holdings, period))


def retirement_calculator(current_age, current_savings, monthly_save,
                           retire_age=65, annual_return=6.0, monthly_expense=50000):
    """退休規劃試算。用法：retirement_calculator <年齡> <現有資產萬元> <月儲蓄元> [退休年齡] [年報酬%] [月支出元]"""
    print(fetch_retirement_calculator(int(current_age), float(current_savings), float(monthly_save),
                                      int(retire_age), float(annual_return), float(monthly_expense)))


def loan_calculator(principal, annual_rate, years, loan_type="等額本息"):
    """貸款試算。用法：loan_calculator <貸款萬元> <年利率%> <年數> [等額本息|等額本金]"""
    print(fetch_loan_calculator(float(principal), float(annual_rate), int(years), loan_type))


def compound_calculator(principal, annual_rate, years, monthly_add=0, compound_freq=12):
    """複利計算器。用法：compound_calculator <本金> <年報酬%> <年數> [每月追加] [複利頻率]"""
    print(fetch_compound_calculator(float(principal), float(annual_rate), int(years),
                                    float(monthly_add), int(compound_freq)))


def asset_allocation(age, risk_level="穩健", goal="退休", investment_horizon=None):
    """資產配置建議。用法：asset_allocation <年齡> [保守|穩健|積極] [目標] [投資年數]"""
    print(fetch_asset_allocation(int(age), risk_level, goal,
                                 int(investment_horizon) if investment_horizon else None))


def tw_tax_calculator(dividend_income, other_income=0, tax_bracket=None, sell_amount=0):
    """台股稅務試算。用法：tw_tax_calculator <股利所得元> [其他收入元] [稅率%] [賣出金額元]"""
    print(fetch_tw_tax_calculator(float(dividend_income), float(other_income),
                                  float(tax_bracket) if tax_bracket else None, float(sell_amount)))


def currency_converter(amount, from_currency, to_currency):
    """外幣換算。用法：currency_converter <金額> <來源幣別> <目標幣別>"""
    print(fetch_currency_converter(float(amount), from_currency, to_currency))


def get_fund(symbol):
    """基金查詢。用法：get_fund <基金代號>"""
    print(fetch_fund(symbol))


def get_reits(symbol):
    """REITs查詢。用法：get_reits <REITs代號>"""
    print(fetch_reits(symbol))


def inflation_adjusted(nominal_return, years, amount, inflation_rate=2.0):
    """通膨調整報酬。用法：inflation_adjusted <名目報酬%> <年數> <本金元> [通膨率%]"""
    print(fetch_inflation_adjusted(float(nominal_return), int(years), float(amount), float(inflation_rate)))


def defi_calculator(principal_usd, apy, days, compound=True, protocol=""):
    """DeFi收益試算。用法：defi_calculator <本金USD> <APY%> <天數> [複利true/false] [協議名]"""
    c = compound if isinstance(compound, bool) else str(compound).lower() != "false"
    print(fetch_defi_calculator(float(principal_usd), float(apy), int(days), c, str(protocol)))


def gold_calculator(weight, unit="公克", currency="TWD"):
    """黃金換算。用法：gold_calculator <重量> [公克|錢|兩|盎司] [TWD|USD]"""
    print(fetch_gold_calculator(float(weight), unit, currency))


def forex_deposit(amount_twd, currency, annual_rate, months, buy_rate=None, sell_rate=None):
    """外幣定存試算。用法：forex_deposit <台幣本金> <幣別> <年利率%> <月數> [買入匯率] [賣出匯率]"""
    print(fetch_forex_deposit(float(amount_twd), currency, float(annual_rate), int(months),
                               float(buy_rate) if buy_rate else None,
                               float(sell_rate) if sell_rate else None))


def financial_health(monthly_income, monthly_expense, total_assets, total_debt,
                     emergency_fund_months=0, has_insurance=False, investment_ratio=0):
    """財務健康診斷。用法：financial_health <月收入> <月支出> <總資產> <總負債> [備用金月數] [有保險y/n] [投資比例%]"""
    ins = has_insurance if isinstance(has_insurance, bool) else str(has_insurance).lower() in ("y", "true", "yes", "1")
    print(fetch_financial_health(float(monthly_income), float(monthly_expense),
                                 float(total_assets), float(total_debt),
                                 float(emergency_fund_months), ins, float(investment_ratio)))


def deep_research(topic, lang="zh-tw", depth=5):
    """深度研究。用法：deep_research <主題> [zh-tw|en] [深度3-8]"""
    print(fetch_deep_research(topic, lang, int(depth)))


def fact_check(claim, lang="zh-tw"):
    """事實查核。用法：fact_check <要查核的說法>"""
    print(fetch_fact_check(claim, lang))


def timeline_events(topic, lang="zh-tw"):
    """時間軸整理。用法：timeline_events <主題>"""
    print(fetch_timeline_events(topic, lang))


def sentiment_scan(topic, lang="zh-tw"):
    """輿情掃描。用法：sentiment_scan <話題>"""
    print(fetch_sentiment_scan(topic, lang))


def compare_analysis(items_str, context=""):
    """多項比較。用法：compare_analysis <A,B,C> [背景說明]"""
    items = [i.strip() for i in items_str.split(",")]
    print(fetch_compare_analysis(items, None, context))


def pros_cons_analysis(subject, context="", lang="zh-tw"):
    """優缺點分析。用法：pros_cons_analysis <主題> [背景] [語言]"""
    print(fetch_pros_cons_analysis(subject, context, lang))


def research_report(topic, purpose="一般研究", lang="zh-tw"):
    """研究報告。用法：research_report <主題> [目的] [語言]"""
    print(fetch_research_report(topic, purpose, lang))


def opinion_writer(topic, stance="中立", style="正式"):
    """觀點撰寫。用法：opinion_writer <主題> [支持|反對|中立|批判] [正式|輕鬆|犀利]"""
    print(fetch_opinion_writer(topic, stance, style))


def trend_forecast(topic, timeframe="全部", lang="zh-tw"):
    """趨勢預測。用法：trend_forecast <主題> [短期|中期|長期|全部]"""
    print(fetch_trend_forecast(topic, timeframe, lang))


def debate_simulator(motion, lang="zh-tw"):
    """辯論模擬。用法：debate_simulator <辯論題目>"""
    print(fetch_debate_simulator(motion, lang))


def academic_search(query, field="", lang="en"):
    """學術論文搜尋。用法：academic_search <關鍵字> [領域] [語言]"""
    print(fetch_academic_search(query, field, lang))


def health_research(topic, lang="zh-tw"):
    """健康資訊搜尋。用法：health_research <主題>"""
    print(fetch_health_research(topic, lang))


def law_research(topic, jurisdiction="台灣", lang="zh-tw"):
    """法規查詢。用法：law_research <主題> [地區] [語言]"""
    print(fetch_law_research(topic, jurisdiction, lang))


def person_research(name, context="", lang="zh-tw"):
    """人物研究。用法：person_research <姓名> [背景說明]"""
    print(fetch_person_research(name, context, lang))


def company_research(company, lang="zh-tw"):
    """公司深度研究。用法：company_research <公司名稱或代號>"""
    print(fetch_company_research(company, lang))


def product_review(product, category="", lang="zh-tw"):
    """產品評測彙整。用法：product_review <產品名稱> [類別]"""
    print(fetch_product_review(product, category, lang))


def travel_research(destination, days=None, style="", lang="zh-tw"):
    """旅遊研究。用法：travel_research <目的地> [天數] [風格]"""
    print(fetch_travel_research(destination, int(days) if days else None, style, lang))


def job_market(job_title, location="台灣", lang="zh-tw"):
    """職涯市場分析。用法：job_market <職位名稱> [地區]"""
    print(fetch_job_market(job_title, location, lang))


def impact_analysis(event, scope_str="", lang="zh-tw"):
    """影響力分析。用法：impact_analysis <事件> [個人,企業,社會,經濟]"""
    scope = [s.strip() for s in scope_str.split(",")] if scope_str else None
    print(fetch_impact_analysis(event, scope, lang))


def scenario_planning(topic, horizon="", lang="zh-tw"):
    """情境規劃。用法：scenario_planning <主題> [時間範圍]"""
    print(fetch_scenario_planning(topic, horizon, lang))


def decision_helper(question, options_str="", criteria_str=""):
    """決策輔助。用法：decision_helper <決策問題> [選項A,選項B] [考量1,考量2]"""
    options = [o.strip() for o in options_str.split(",")] if options_str else None
    criteria = [c.strip() for c in criteria_str.split(",")] if criteria_str else None
    print(fetch_decision_helper(question, options, criteria))


def devil_advocate(position, lang="zh-tw"):
    """魔鬼代言人。用法：devil_advocate <要被挑戰的觀點>"""
    print(fetch_devil_advocate(position, lang))


def summary_writer(topic, max_points=7, lang="zh-tw"):
    """多來源摘要。用法：summary_writer <主題> [重點數] [語言]"""
    print(fetch_summary_writer(topic, int(max_points), lang))


def key_insights(topic, count=5, lang="zh-tw"):
    """洞察萃取。用法：key_insights <主題> [數量] [語言]"""
    print(fetch_key_insights(topic, int(count), lang))


def bias_detector(topic, lang="zh-tw"):
    """偏見偵測。用法：bias_detector <議題>"""
    print(fetch_bias_detector(topic, lang))


def second_opinion(question, experts_str="", lang="zh-tw"):
    """多專家視角。用法：second_opinion <問題> [專家A,專家B,...]"""
    experts = [e.strip() for e in experts_str.split(",")] if experts_str else None
    print(fetch_second_opinion(question, experts, lang))


def brainstorm(problem, count=8, style="實用", lang="zh-tw"):
    """腦力激盪。用法：brainstorm <問題> [數量] [實用|創意|顛覆]"""
    print(fetch_brainstorm(problem, int(count), style, lang))


def benchmark_analysis(subject, industry="", lang="zh-tw"):
    """標竿分析。用法：benchmark_analysis <對象> [產業]"""
    print(fetch_benchmark_analysis(subject, industry, lang))


def steel_man(opposing_view, own_position="", lang="zh-tw"):
    """鋼人論證。用法：steel_man <對立觀點> [自己的立場]"""
    print(fetch_steel_man(opposing_view, own_position, lang))


def socratic_questioning(topic, depth=5, lang="zh-tw"):
    """蘇格拉底式提問。用法：socratic_questioning <主題> [層數]"""
    print(fetch_socratic_questioning(topic, int(depth), lang))


def analogy_maker(concept, audience="一般大眾", count=3, lang="zh-tw"):
    """類比說明。用法：analogy_maker <概念> [受眾] [數量]"""
    print(fetch_analogy_maker(concept, audience, int(count), lang))


def narrative_builder(topic, key_message="", audience="", lang="zh-tw"):
    """敘事架構。用法：narrative_builder <主題> [核心訊息] [受眾]"""
    print(fetch_narrative_builder(topic, key_message, audience, lang))


def critique_writer(subject, type_="觀點", lang="zh-tw"):
    """批判性評析。用法：critique_writer <對象> [文章|政策|作品|計劃|觀點]"""
    print(fetch_critique_writer(subject, type_, lang))


def position_statement(issue, stance, lang="zh-tw"):
    """立場聲明。用法：position_statement <議題> <支持|反對|有條件支持>"""
    print(fetch_position_statement(issue, stance, lang))


def ocr_click(target_text, monitor=1, click_type="click"):
    """OCR找字點擊。用法：ocr_click <目標文字> [螢幕1/2/3] [click|double_click|right_click]"""
    print(fetch_ocr_click(target_text, int(monitor), click_type))


def vision_locate(description, monitor=1, action="click"):
    """視覺定位點擊。用法：vision_locate <描述> [螢幕] [click|double_click|locate_only]"""
    print(fetch_vision_locate(description, int(monitor), action))


def screen_workflow(steps_json):
    """螢幕工作流。用法：screen_workflow '<JSON步驟陣列>'"""
    import json
    steps = json.loads(steps_json) if isinstance(steps_json, str) else steps_json
    print(fetch_screen_workflow(steps))


def app_navigator(app, task, input_text="", monitor=1):
    """App導航。用法：app_navigator <App名> <任務描述> [輸入文字] [螢幕]"""
    print(fetch_app_navigator(app, task, input_text, int(monitor)))


def wait_and_click(target_text, timeout=15, monitor=1, action_after="click"):
    """等待出現後點擊。用法：wait_and_click <目標文字> [超時秒數] [螢幕] [click|none]"""
    print(fetch_wait_and_click(target_text, int(timeout), int(monitor), action_after))


def drag_drop(from_x=None, from_y=None, to_x=None, to_y=None, from_text="", to_text="", monitor=1, duration=0.5):
    """拖曳操作。用法：drag_drop <from_x> <from_y> <to_x> <to_y> 或 drag_drop from_text=<文字> to_text=<文字>"""
    print(fetch_drag_drop(
        int(from_x) if from_x else None, int(from_y) if from_y else None,
        int(to_x) if to_x else None, int(to_y) if to_y else None,
        str(from_text), str(to_text), int(monitor), float(duration)))


# ── 新增 116 個缺失工具（同步 bot.py）──────────────────

def analyze_pdf_tool(path, max_chars=4000):
    """分析 PDF。用法：analyze_pdf <路徑> [最大字數]"""
    try:
        import pdfplumber
        with pdfplumber.open(path) as pdf:
            total_pages = len(pdf.pages)
            texts = []
            for page in pdf.pages[:20]:
                t = page.extract_text()
                if t:
                    texts.append(t)
        full_text = "\n".join(texts)
        if not full_text.strip():
            print("PDF 無法提取文字（可能是掃描圖片 PDF）"); return
        result = f"📄 PDF 分析（共 {total_pages} 頁）\n\n{full_text[:int(max_chars)]}"
        if len(full_text) > int(max_chars):
            result += f"\n\n（內容已截斷，共 {len(full_text)} 字）"
        print(result)
    except Exception as e:
        print(f"PDF 讀取失敗：{e}")


def audio_process_tool(action, input_path, output="", start_ms=0, end_ms=0):
    """音訊處理。用法：audio_process <convert|trim|info> <輸入> [輸出] [起始ms] [結束ms]"""
    print(execute_audio_process(action, input_path, output, int(start_ms), int(end_ms)))


def auto_skill_tool(action, goal="", skill_name="", code=""):
    """自動技能。用法：auto_skill <create|list|run|delete> [目標] [名稱] [程式碼]"""
    print(execute_auto_skill(action, goal, skill_name, code))


def auto_trade_tool(action, symbol="", amount=0.0, price=0.0, order_type="market"):
    """自動交易。用法：auto_trade <add|remove|list|simulate> [代號] [數量] [價格] [類型]"""
    print(execute_auto_trade(action, symbol, float(amount), float(price), order_type))


def automation_tool(action, condition_type="", condition_value="", command="", duration=60.0, layout="side_by_side", x=0, y=0, w=0, h=0, keyword="", output=""):
    """自動化。用法：automation <if_then|window_arrange|region_ocr|...> [參數...]"""
    print(execute_automation(action, condition_type, condition_value, command, float(duration), layout, int(x), int(y), int(w), int(h), keyword, output))


def barcode_tool(image_path=""):
    """掃描條碼/QR Code。用法：barcode [圖片路徑]"""
    print(execute_barcode(image_path))


def bluetooth_tool(action, mac=""):
    """藍牙操作。用法：bluetooth <scan|connect|disconnect|list> [MAC]"""
    print(execute_bluetooth(action, mac))


def browser_advanced_tool(action, selector="", value="", name="", tab_index=0, timeout=30.0, url_pattern=""):
    """進階瀏覽器。用法：browser_advanced <action> [selector] [value] [name] [tab] [timeout] [url_pattern]"""
    print(execute_browser_advanced(action, selector, value, name, int(tab_index), float(timeout), url_pattern))


def browser_control_tool(action, url="", selector="", text=""):
    """瀏覽器控制。用法：browser_control <open|goto|click|type|get_text|screenshot|close> [url] [selector] [text]"""
    print(execute_browser_control(action, url, selector, text))


def calendar_tool(action, days=7, title="", start="", end="", description=""):
    """Google 日曆。用法：calendar <list|add> [天數] [標題] [開始] [結束] [說明]"""
    print(execute_calendar(action, int(days), title, start, end, description))


def clipboard_tool(action, text=""):
    """剪貼簿。用法：clipboard <get|set|history> [文字]"""
    print(execute_clipboard(action, text))


def clipboard_image_tool(action, path=""):
    """剪貼簿圖片。用法：clipboard_image <get|set> [路徑]"""
    print(execute_clipboard_image(action, path))


def cloud_storage_tool(action, path, drive_id="root"):
    """雲端儲存。用法：cloud_storage <upload|download|list> <路徑> [drive_id]"""
    print(execute_cloud_storage(action, path, drive_id))


def compare_stocks_tool(symbols_str, metrics_str="all"):
    """比較股票。用法：compare_stocks <代號1,代號2,...> [指標]"""
    _compare_stocks = compare_stocks
    symbols = symbols_str.split(",")
    metrics = metrics_str.split(",") if metrics_str != "all" else None
    print(_compare_stocks(symbols, metrics))


def database_tool(type_, db, sql, name=""):
    """資料庫。用法：database <sqlite|mysql> <路徑/host> <SQL> [名稱]"""
    print(execute_database(type_, db, sql, name))


def ddg_search_tool(query, region="zh-tw", max_results=5):
    """DuckDuckGo 搜尋。用法：ddg_search <關鍵字> [地區] [數量]"""
    print(execute_ddg_search(query, region, int(max_results)))


def desktop_control_tool(action, x=None, y=None, text=None, app=None, direction="down", amount=3, monitor=None):
    """桌面控制。用法：desktop_control <action> [x] [y] [text] [app] [direction] [amount] [monitor]"""
    result = execute_desktop_control(action, x=int(x) if x else None, y=int(y) if y else None,
                                     text=text, app=app, direction=direction, amount=int(amount),
                                     monitor=int(monitor) if monitor else None)
    print(result.get("message", str(result)))


def device_manager_tool(action, name="", keyword=""):
    """裝置管理員。用法：device_manager <list|enable|disable> [名稱] [關鍵字]"""
    print(execute_device_manager(action, name, keyword))


def disk_backup_tool(action, src="", dest=""):
    """磁碟備份。用法：disk_backup <backup|restore|list> [來源] [目標]"""
    print(execute_disk_backup(action, src, dest))


def display_tool(action, level=None):
    """顯示設定。用法：display <get_brightness|set_brightness|get_resolution|list_resolutions> [亮度]"""
    print(execute_display(action, int(level) if level else None))


def docker_tool(action, name=""):
    """Docker 操作。用法：docker <list|start|stop|logs|images> [容器名]"""
    print(execute_docker(action, name))


def document_control_tool(action, path, content="", sheet=None):
    """文件控制。用法：document_control <read|write|list_sheets> <路徑> [內容] [工作表]"""
    print(execute_document(action, path, content, sheet))


def download_file_tool(url, save_path=""):
    """下載檔案。用法：download_file <URL> [儲存路徑]"""
    print(execute_download_file(url, save_path))


def dropbox_tool(action, local, remote, token=""):
    """Dropbox 操作。用法：dropbox <upload|download> <本地> <遠端> [token]"""
    print(execute_dropbox(action, local, remote, token))


def email_control_tool(host, user, password, folder="INBOX", count=5):
    """Email 控制（讀取）。用法：email_control <host> <user> <password> [folder] [count]"""
    print(execute_email_read(host, user, password, folder, int(count)))


def emotion_detect_tool(action, text="", image_path=""):
    """情緒偵測。用法：emotion_detect <text|image|both> [文字] [圖片路徑]"""
    print(execute_emotion_detect(action, text, image_path))


def encrypt_file_tool(action, path, password):
    """加解密檔案。用法：encrypt_file <encrypt|decrypt> <路徑> <密碼>"""
    print(execute_encrypt_file(action, path, password))


def env_var_tool(action, name="", value="", permanent="false"):
    """環境變數。用法：env_var <get|set|list|delete> [名稱] [值] [permanent]"""
    print(execute_env_var(action, name, value, permanent.lower() == "true"))


def file_system_tool(action, path="", dest="", content="", keyword=""):
    """檔案系統。用法：file_system <list|read|write|delete|copy|move|search|info> [路徑] [目標] [內容] [關鍵字]"""
    print(execute_file_system(action, path, dest, content, keyword))


def file_tools_tool(action, path, dest="", pattern="", replacement="", ext=""):
    """檔案工具。用法：file_tools <batch_rename|sync> <路徑> [目標] [pattern] [replacement] [ext]"""
    print(execute_file_tools(action, path, dest, pattern, replacement, ext))


def file_transfer_tool(action, source, dest=""):
    """檔案傳輸。用法：file_transfer <zip|unzip|download|upload> <來源> [目標]"""
    print(execute_file_transfer(action, source, dest))


def find_image_on_screen_tool(image_path, confidence=0.8):
    """在螢幕上找圖。用法：find_image_on_screen <圖片路徑> [信心值]"""
    print(execute_find_image(image_path, float(confidence)))


def generate_image_tool(prompt, width=512, height=512, overlay_text=""):
    """生成圖片。用法：generate_image <prompt> [寬] [高] [疊加文字]"""
    image_bytes = fetch_image(prompt, int(width), int(height))
    if image_bytes and overlay_text:
        image_bytes = add_text_overlay(image_bytes, overlay_text)
    if image_bytes:
        out = str(Path.home() / "Desktop" / f"gen_{int(time.time())}.png")
        with open(out, "wb") as f:
            f.write(image_bytes)
        print(f"✅ 圖片已生成：{out}")
    else:
        print("圖片生成失敗")


def get_candlestick_chart_tool(symbol, period="3mo"):
    """K線圖。用法：get_candlestick_chart <代號> [期間]"""
    img_bytes, pattern_str = generate_candlestick(symbol, period)
    if img_bytes:
        out = str(Path.home() / "Desktop" / f"candlestick_{symbol}_{int(time.time())}.png")
        with open(out, "wb") as f:
            f.write(img_bytes)
        print(f"✅ K線圖已儲存：{out}\n\n{pattern_str}")
    else:
        print(pattern_str)


def get_crypto_tool(coin, vs_currency="usd"):
    """查詢加密貨幣。用法：get_crypto <幣種> [vs貨幣]"""
    print(fetch_crypto(coin, vs_currency))


def get_earnings_tool(symbol):
    """查詢財報。用法：get_earnings <代號>"""
    print(fetch_earnings(symbol))


def get_etf_tool(symbol):
    """查詢 ETF。用法：get_etf <代號>"""
    print(fetch_etf(symbol))


def get_finance_news_tool(source="all", count=5):
    """財經新聞。用法：get_finance_news [來源] [數量]"""
    print(fetch_finance_news(source, int(count)))


def get_forex_tool(pair):
    """查詢匯率。用法：get_forex <貨幣對>"""
    print(fetch_forex(pair))


def get_fundamentals_tool(symbol):
    """查詢基本面。用法：get_fundamentals <代號>"""
    print(fetch_fundamentals(symbol))


def get_macro_tool(indicator):
    """查詢總經指標。用法：get_macro <cpi|gdp|unemployment|fed_rate|nonfarm>"""
    print(fetch_macro(indicator))


def get_market_sentiment_tool():
    """查詢市場情緒。用法：get_market_sentiment"""
    print(fetch_market_sentiment())


def get_stock_tool(symbol, period="1mo"):
    """查詢股票。用法：get_stock <代號> [期間]"""
    print(fetch_stock(symbol, period))


def get_stock_advanced_tool(symbol, indicators="macd,bb,kd"):
    """進階技術分析。用法：get_stock_advanced <代號> [指標]"""
    ind_list = indicators.split(",") if indicators else None
    print(fetch_stock_advanced(symbol, ind_list))


def get_weather_tool(city):
    """查詢天氣。用法：get_weather <城市>"""
    print(fetch_weather(city))


def git_tool(action, repo=".", message="", branch="master"):
    """Git 操作。用法：git <status|log|pull|add|commit|push|diff> [repo] [message] [branch]"""
    print(execute_git(action, repo, message, branch))


def goal_manager_tool(action, goal="", goal_id="", steps="", priority="normal"):
    """目標管理。用法：goal_manager <add|list|update|delete|progress> [目標] [id] [步驟] [優先]"""
    print(execute_goal_manager(action, goal, goal_id, steps, priority))


def google_trends_tool(keywords_str, timeframe="today 3-m", geo="TW"):
    """Google Trends。用法：google_trends <關鍵字1,關鍵字2,...> [時間] [地區]"""
    keywords = keywords_str.split(",")
    print(fetch_google_trends(keywords, timeframe, geo))


def hardware_tool():
    """硬體監控。用法：hardware"""
    print(execute_hardware())


def image_edit_tool(action, path, *params):
    """圖片編輯。用法：image_edit <crop|resize|text|merge|rotate|flip> <路徑> [參數...]"""
    print(execute_image_edit(action, path, *params))


def image_tools_tool(action, path="", quality=75, width=0, height=0, target_lang="zh-TW"):
    """圖片工具。用法：image_tools <compress|batch|ocr_translate> [路徑] [quality] [width] [height] [lang]"""
    print(execute_image_tools(action, path, int(quality), int(width), int(height), target_lang))


def knowledge_base_tool(action, content="", query="", tag="", kb_id=""):
    """知識庫。用法：knowledge_base <add|search|list|delete|export> [content] [query] [tag] [id]"""
    print(execute_knowledge_base(action, content, query, tag, kb_id))


def long_term_memory_tool(action, chat_id="0", content="", memory_id=""):
    """長期記憶。用法：long_term_memory <save|list|delete> [chat_id] [內容] [id]"""
    if action == "save":
        memory_save(int(chat_id), content)
    elif action == "list":
        memory_list(int(chat_id))
    elif action == "delete":
        memory_del(int(chat_id), int(memory_id))
    else:
        print(f"可用操作：save, list, delete")


def lookup_tool(action, ip="", amount=1.0, from_cur="USD", to_cur="TWD"):
    """查詢工具。用法：lookup <ip|currency> [ip] [amount] [from] [to]"""
    print(execute_lookup(action, ip, float(amount), from_cur, to_cur))


def manage_schedule_tool(action, name="", time_str="", script=""):
    """排程管理。用法：manage_schedule <list|add|remove|enable|disable> [名稱] [時間] [腳本]"""
    print(execute_manage_schedule(action, name, time_str, script))


def media_tool(action, device_name=""):
    """媒體控制。用法：media <play_pause|next|prev|stop|volume_up|volume_down|mute|list_devices|switch> [裝置名]"""
    print(execute_media(action, device_name))


def monitor_config_tool():
    """螢幕設定。用法：monitor_config"""
    print(execute_monitor_config())


def multi_deploy_tool(action, remote_host="", remote_user="", remote_pass="", remote_path="/tmp/niu_bot"):
    """多機部署。用法：multi_deploy <deploy|status|rollback> [host] [user] [pass] [path]"""
    print(execute_multi_deploy(action, remote_host, remote_user, remote_pass, remote_path))


def multi_perspective_tool(topic, lang="zh-tw"):
    """多角度分析。用法：multi_perspective <主題> [語言]"""
    print(multi_perspective(topic, lang))


def network_config_tool(action, name="", ip="", dns1="", dns2="", domain="", duration=10):
    """網路設定。用法：network_config <list|set_ip|set_dns|get_dns|hosts_list|hosts_add|hosts_remove|traffic> [參數...]"""
    print(execute_network_config(action, name, ip, dns1, dns2, domain, int(duration)))


def network_diag_tool(action, host, ports="22,80,443,3306,3389,8080"):
    """網路診斷。用法：network_diag <ping|traceroute|portscan> <host> [ports]"""
    print(execute_network_diag(action, host, ports))


def news_monitor_tool(action, keywords="", interval=300, duration=3600):
    """新聞監控。用法：news_monitor <start|stop|status> [關鍵字] [間隔秒] [持續秒]"""
    print(execute_news_monitor(action, keywords, int(interval), int(duration)))


def news_search_tool(query, lang="zh-TW", count=6):
    """新聞搜尋。用法：news_search <關鍵字> [語言] [數量]"""
    print(search_news(query, lang, int(count)))


def nlp_tool(action, text):
    """NLP 工具。用法：nlp <summarize|sentiment|keywords|ner> <文字>"""
    print(execute_nlp(action, text))


def osint_search_tool(action, query="", target="", limit=10):
    """OSINT 搜尋。用法：osint_search <web_search|news_search|ip_osint|domain_osint|top_news> [query] [target] [limit]"""
    print(execute_osint_search(action, query, target, int(limit)))


def password_mgr_tool(action, site, master, username="", password=""):
    """密碼管理。用法：password_mgr <save|get|list|delete> <網站> <主密碼> [帳號] [密碼]"""
    print(execute_password_mgr(action, site, master, username, password))


def pdf_edit_tool(action, path="", output="", paths="", text=""):
    """PDF 編輯。用法：pdf_edit <merge|split|watermark|info> [路徑] [輸出] [paths] [text]"""
    print(execute_pdf_edit(action, path, output, paths, text))


def pdf_image_tool(path, output_dir="", dpi=150):
    """PDF 轉圖片。用法：pdf_image <路徑> [輸出資料夾] [dpi]"""
    print(execute_pdf_to_image(path, output_dir, int(dpi)))


def pentest_tool(action, target="", port_range="1-1000", timeout=2):
    """滲透測試。用法：pentest <scan|vuln_check|banner_grab|ssl_check> [target] [port_range] [timeout]"""
    print(execute_pentest(action, target, port_range, int(timeout)))


def portfolio_tool(action, chat_id=0, symbol="", shares=0, cost=0):
    """投資組合。用法：portfolio <add|remove|view|clear> [chat_id] [symbol] [shares] [cost]"""
    print(execute_portfolio(action, int(chat_id), symbol, float(shares), float(cost)))


def power_control_tool(action):
    """電源控制。用法：power_control <sleep|restart|shutdown|hibernate|lock>"""
    print(execute_power(action))


def pptx_control_tool(action, path, slides=""):
    """PowerPoint 控制。用法：pptx_control <read|create|add_slide> <路徑> [slides_json]"""
    print(execute_pptx(action, path, slides))


def proactive_alert_tool(action, name="", condition="", threshold="", target="", interval=60):
    """主動警報。用法：proactive_alert <add|remove|list|status> [名稱] [條件] [閾值] [目標] [間隔秒]"""
    print(execute_proactive_alert(action, name, condition, threshold, target, int(interval)))


def ptt_search_tool(keyword, board="Gossiping", count=5):
    """PTT 搜尋。用法：ptt_search <關鍵字> [看板] [數量]"""
    _ptt_search = ptt_search
    print(_ptt_search(keyword, board, int(count)))


def push_notify_tool(platform, message, webhook_or_token):
    """推播通知。用法：push_notify <discord|line|slack> <訊息> <webhook_or_token>"""
    print(execute_push_notify(platform, message, webhook_or_token))


def qr_code_tool(action, content="", path="", duration=30.0):
    """QR Code。用法：qr_code <generate|scan|watch> [content] [path] [duration]"""
    print(execute_qr_code(action, content, path, float(duration)))


def read_screen_tool(question="描述螢幕上有什麼", monitor=1):
    """讀取螢幕。用法：read_screen [問題] [螢幕號]"""
    print(fetch_read_screen(question, int(monitor)))


def read_webpage_tool(url, max_chars=3000):
    """讀取網頁。用法：read_webpage <URL> [最大字數]"""
    print(read_webpage(url, int(max_chars)))


def registry_tool(action, key, value_name="", value=""):
    """登錄檔操作。用法：registry <read|write|delete|list> <key> [value_name] [value]"""
    print(execute_registry(action, key, value_name, value))


def reminder_tool(time_str, message):
    """提醒。用法：reminder <HH:MM或秒數> <訊息>"""
    print(execute_reminder(time_str, message))


def report_tool(title, data_json, output=""):
    """報告生成。用法：report <標題> <資料JSON> [輸出路徑]"""
    print(execute_report(title, data_json, output))


def restore_point_tool(action, description=""):
    """系統還原點。用法：restore_point <create|list> [說明]"""
    print(execute_restore_point(action, description))


def run_code_tool(type_, code):
    """執行程式碼。用法：run_code <python|powershell|cmd|javascript> <code>"""
    print(execute_run_code(type_, code))


def screen_vision_tool(question="請描述這個畫面上有什麼，以及目前電腦在做什麼事。"):
    """螢幕視覺分析。用法：screen_vision [問題]"""
    analysis, img_bytes = execute_screen_vision(question)
    if img_bytes:
        out = str(Path.home() / "Desktop" / f"vision_{int(time.time())}.png")
        with open(out, "wb") as f:
            f.write(img_bytes)
        print(f"{analysis}\n\n截圖已儲存：{out}")
    else:
        print(analysis)


def scroll_at_tool(direction="down", amount=3, x=None, y=None, monitor=1, description=""):
    """指定位置滾動。用法：scroll_at [方向] [格數] [x] [y] [螢幕號] [描述]"""
    print(fetch_scroll_at(direction, int(amount),
                          int(x) if x and x != "None" else None,
                          int(y) if y and y != "None" else None,
                          int(monitor), description))


def self_benchmark_tool(action):
    """自我評測。用法：self_benchmark <run|report|compare>"""
    print(execute_self_benchmark(action))


def send_email_tool(to, subject, body):
    """發送 Email。用法：send_email <收件人> <主旨> <內容>"""
    print(execute_send_email(to, subject, body))


def send_voice_tool(text, voice="zh-CN-YunxiNeural"):
    """生成語音。用法：send_voice <文字> [語音]"""
    try:
        ogg_data = generate_voice_ogg(text, voice)
        out = str(Path.home() / "Desktop" / f"voice_{int(time.time())}.ogg")
        with open(out, "wb") as f:
            f.write(ogg_data)
        print(f"✅ 語音已生成：{out}")
    except Exception as e:
        print(f"語音生成失敗：{e}")


def smart_home_tool(action, device="", value="", host="", token=""):
    """智慧家居。用法：smart_home <list|control|status|scene> [device] [value] [host] [token]"""
    print(execute_smart_home(action, device, value, host, token))


def software_tool(action, name="", keyword=""):
    """軟體管理。用法：software <list|install|uninstall> [名稱] [關鍵字]"""
    print(execute_software(action, name, keyword))


def ssh_sftp_tool(action, host, user, password, command="", local="", remote="", port=22):
    """SSH/SFTP。用法：ssh_sftp <run|upload|download> <host> <user> <pass> [command] [local] [remote] [port]"""
    print(execute_ssh_sftp(action, host, user, password, command, local, remote, int(port)))


def startup_tool(action, name="", command=""):
    """開機自啟。用法：startup <list|add|remove> [名稱] [指令]"""
    print(execute_startup(action, name, command))


def system_monitor_tool(action, target=""):
    """系統監控。用法：system_monitor <sysinfo|process_list|process_kill|disk_usage|network|battery> [target]"""
    print(execute_system_monitor(action, target))


def system_tools_tool(action, **kwargs):
    """系統工具。用法：system_tools <event_log|usb_list|firewall_list|printer_list|...> [參數...]"""
    print(execute_system_tools(action, **kwargs))


def think_as_tool(person, question, list_available="false"):
    """角色思考。用法：think_as <人物> <問題> [list]"""
    print(execute_think_as(person, question, list_available.lower() == "true"))


def threat_intel_tool(action, target="", api_key=""):
    """威脅情報。用法：threat_intel <ip_check|domain_check|hash_check|cve_search> [target] [api_key]"""
    print(execute_threat_intel(action, target, api_key))


def todo_list_tool(action, task="", todo_id=0):
    """任務清單。用法：todo_list <add|list|done|delete|clear> [task] [id]"""
    print(execute_todo(action, task, int(todo_id)))


def tts_advanced_tool(action, text="", voice="zh-CN-YunxiNeural"):
    """進階 TTS。用法：tts_advanced <speak|list_voices> [文字] [語音]"""
    print(execute_tts_advanced(action, text, voice))


def user_account_tool(action, username="", password=""):
    """使用者帳戶。用法：user_account <list|create|delete> [帳號] [密碼]"""
    print(execute_user_account(action, username, password))


def video_gif_tool(path, start=0, duration=5.0, output="", fps=10):
    """影片轉 GIF。用法：video_gif <路徑> [起始秒] [持續秒] [輸出] [fps]"""
    print(execute_video_gif(path, float(start), float(duration), output, int(fps)))


def video_process_tool(action, path, second=0, start=0, end=0, output=""):
    """影片處理。用法：video_process <screenshot|trim|info|to_gif> <路徑> [秒數] [起始] [結束] [輸出]"""
    print(execute_video_process(action, path, float(second), float(start), float(end), output))


def virtual_desktop_tool(action):
    """虛擬桌面。用法：virtual_desktop <left|right|new>"""
    print(execute_vdesktop(action))


def voice_cmd_tool(action, duration=300.0, language="zh-TW"):
    """語音命令。用法：voice_cmd <start|stop|status> [持續秒] [語言]"""
    print(execute_voice_cmd(action, float(duration), language))


def voice_id_tool(action, name="", audio_path="", duration=5):
    """聲紋辨識。用法：voice_id <register|identify|list|delete> [名稱] [音檔] [秒數]"""
    print(execute_voice_id(action, name, audio_path, int(duration)))


def volume_tool(action, level=None):
    """音量控制。用法：volume <get|set|mute|unmute> [音量]"""
    print(execute_volume(action, int(level) if level else None))


def vpn_tool(action, name="", user="", password=""):
    """VPN 控制。用法：vpn <list|connect|disconnect> [名稱] [帳號] [密碼]"""
    print(execute_vpn(action, name, user, password))


def wait_seconds_tool(seconds):
    """等待。用法：wait_seconds <秒數>"""
    print(execute_wait_seconds(float(seconds)))


def web_scrape_tool(action, url="", selector="body", interval=2.0, region="full"):
    """網頁爬取。用法：web_scrape <scrape|monitor|screenshot> [url] [selector] [interval] [region]"""
    print(execute_web_scrape(action, url, selector, float(interval), region))


def webpage_shot_tool(action, url, selector="body", interval=60.0, duration=3600.0):
    """網頁截圖。用法：webpage_shot <screenshot|monitor> <url> [selector] [interval] [duration]"""
    print(execute_webpage_shot(action, url, selector, float(interval), float(duration)))


def wikipedia_search_tool(query, lang="zh"):
    """Wikipedia 搜尋。用法：wikipedia_search <關鍵字> [語言]"""
    _wikipedia_search = wikipedia_search
    print(_wikipedia_search(query, lang))


def win_notify_relay_tool(action, duration=3600.0, filter_app=""):
    """Windows 通知轉發。用法：win_notify_relay <start|stop|status> [持續秒] [filter_app]"""
    print(execute_win_notify_relay(action, float(duration), filter_app))


def window_control_tool(action, keyword=""):
    """視窗控制。用法：window_control <list|focus|maximize|minimize|close|restore> [關鍵字]"""
    print(execute_window_control(action, keyword))


def window_manager_tool(action="list", window_name=""):
    """視窗管理員。用法：window_manager <list|focus|maximize|minimize|close> [視窗名]"""
    print(fetch_window_manager(action, window_name))


def windows_update_tool(action):
    """Windows 更新。用法：windows_update <list|install|check>"""
    print(execute_win_update(action))


def workflow_tool(action, name="", steps=""):
    """工作流程。用法：workflow <run|save|list|delete> [名稱] [steps_json]"""
    print(execute_workflow(action, name, steps))


def youtube_summary_tool(url):
    """YouTube 摘要。用法：youtube_summary <URL>"""
    _youtube_summary = youtube_summary
    print(_youtube_summary(url))


def tg_auto_reply_tool(action="start", contact="", stop_time="", duration="30"):
    """Telegram 自動回覆。用法：tg_auto_reply [start|stop] [contact_name] [stop_time]"""
    if action == "stop":
        try:
            with open("C:/Users/blue_/Desktop/測試檔案/.stop_auto_reply", "w") as f:
                f.write("stop")
            print("自動回覆已停止")
        except Exception as e:
            print(f"停止失敗：{e}")
        return
    if not contact:
        print("請提供好友名稱。用法：tg_auto_reply start <好友名稱> <HH:MM>")
        return
    if not stop_time:
        from datetime import datetime, timedelta
        stop_time = (dt.datetime.now() + timedelta(minutes=int(duration))).strftime("%H:%M")
    script = "C:/Users/blue_/claude-telegram-bot/scripts/tg_auto_chat.py"
    subprocess.Popen([sys.executable, script, contact, stop_time], cwd="C:/Users/blue_/claude-telegram-bot")
    print(f"自動回覆已開啟：對象 {contact}，監控到 {stop_time}")


# ── LINE 工具 ──────────────────────────────────────────────────────────

def execute_line_send_msg(contact_name: str, message: str) -> str:
    """LINE 搜尋好友並發送訊息（subprocess 呼叫 line_send_msg.py）"""
    if not contact_name or not message:
        return "請提供好友名稱和訊息內容"

    # GPU 互斥檢查
    available, reason = _check_gpu_available("line_send_msg")
    if not available:
        return reason

    try:
        _set_gpu_active("line_send_msg")
        script = "C:/Users/blue_/claude-telegram-bot/scripts/line_send_msg.py"
        proc = subprocess.Popen(
            [sys.executable, script, contact_name, message],
            cwd="C:/Users/blue_/claude-telegram-bot",
            stdout=subprocess.PIPE, stderr=subprocess.STDOUT,
            encoding="utf-8", errors="replace",
        )
        output, _ = proc.communicate(timeout=120)
        if proc.returncode == 0:
            return f"LINE 訊息已發送給 {contact_name}：{message}"
        else:
            return f"LINE 發送失敗：{output[-500:]}"
    except subprocess.TimeoutExpired:
        proc.kill()
        return "LINE 發送超時（120秒）"
    except Exception as e:
        return f"LINE 發送失敗：{e}"
    finally:
        _release_gpu("line_send_msg")


def line_send_msg_tool(contact="", message=""):
    """LINE 搜尋好友並發送訊息。用法：line_send_msg <好友名稱> <訊息>"""
    if not contact or not message:
        print("請提供好友名稱和訊息。用法：line_send_msg <好友名稱> <訊息>")
        return
    script = "C:/Users/blue_/claude-telegram-bot/scripts/line_send_msg.py"
    proc = subprocess.Popen(
        [sys.executable, script, contact, message],
        cwd="C:/Users/blue_/claude-telegram-bot",
    )
    proc.wait(timeout=120)
    if proc.returncode == 0:
        print(f"LINE 訊息已發送給 {contact}")
    else:
        print(f"LINE 發送失敗 (exit code={proc.returncode})")


# ── 從 bot.py 同步的 29 個 execute_* 函數 ──────────────────────────────

_interval_schedules = {}
_screen_live_running = False


def execute_ai_plan(goal: str) -> str:
    import anthropic, json, time as _time
    _client = anthropic.Anthropic(api_key=os.getenv("ANTHROPIC_API_KEY"))
    response = _client.messages.create(
        model="claude-sonnet-4-6",
        max_tokens=2048,
        system="""你是電腦自動化規劃師。把目標拆解成可執行步驟，以 JSON 陣列回傳：
[{"tool":"click","args":[x,y],"delay":0},{"tool":"type","args":["文字"]}]
可用工具：click/type/press/hotkey/open/screenshot/wait/notify/move/scroll
只回傳 JSON。""",
        messages=[{"role": "user", "content": f"目標：{goal}"}]
    )
    plan_text = response.content[0].text.strip()
    try:
        steps = json.loads(plan_text)
        results = []
        tool_map = {
            "click": lambda a: pyautogui.click(int(a[0]), int(a[1])),
            "type": lambda a: (
                __import__("pyperclip").copy(" ".join(str(x) for x in a)),
                __import__("time").sleep(0.2),
                pyautogui.hotkey("ctrl", "v")
            ),
            "press": lambda a: pyautogui.press(a[0]),
            "hotkey": lambda a: pyautogui.hotkey(*a),
            "open": lambda a: subprocess.Popen(" ".join(str(x) for x in a), shell=True),
            "wait": lambda a: _time.sleep(float(a[0])),
            "move": lambda a: pyautogui.moveTo(int(a[0]), int(a[1]), duration=0.3),
            "scroll": lambda a: pyautogui.scroll(int(a[1]) if a[0]=="up" else -int(a[1])),
        }
        for i, step in enumerate(steps, 1):
            t = step.get("tool"); a = step.get("args", []); d = step.get("delay", 0)
            if d: _time.sleep(d)
            if t in tool_map:
                tool_map[t](a)
                results.append(f"步驟 {i} ✅ {t}")
            else:
                results.append(f"步驟 {i} ⚠️ 未知：{t}")
        return f"目標「{goal}」執行完畢\n" + "\n".join(results)
    except Exception as e:
        return f"規劃結果：{plan_text}\n執行錯誤：{e}"


def execute_api_call(method, url, headers="{}", body="{}"):
    try:
        import json
        h = json.loads(headers) if headers else {}
        b = json.loads(body) if body else None
        resp = requests.request(method.upper(), url, headers=h, json=b, timeout=30)
        try:
            return json.dumps(resp.json(), ensure_ascii=False, indent=2)[:2000]
        except Exception:
            return resp.text[:2000]
    except Exception as e:
        return f"❌ API 呼叫失敗：{e}"


def execute_chrome_bookmarks():
    try:
        import json
        bookmark_path = Path.home() / "AppData/Local/Google/Chrome/User Data/Default/Bookmarks"
        if not bookmark_path.exists():
            return "❌ 找不到 Chrome 書籤（Chrome 未安裝或路徑不同）"
        data = json.loads(bookmark_path.read_text(encoding="utf-8"))
        lines = []
        def _collect(node, indent=0):
            if node.get("type") == "url":
                lines.append("  " * indent + f"🔗 {node['name']}  {node['url']}")
            elif node.get("type") == "folder":
                lines.append("  " * indent + f"📁 {node['name']}")
                for child in node.get("children", []):
                    _collect(child, indent + 1)
        for root in data["roots"].values():
            _collect(root)
        return "📚 Chrome 書籤：\n" + "\n".join(lines[:80])
    except Exception as e:
        return f"❌ 讀取書籤失敗：{e}"


def execute_defender(action, path=""):
    try:
        if action == "status":
            r = subprocess.run(["powershell", "-Command",
                "Get-MpComputerStatus | Select-Object AMRunningMode,RealTimeProtectionEnabled,AntivirusSignatureLastUpdated | Format-List"],
                capture_output=True, text=True, encoding="utf-8", errors="replace")
            return f"🛡️ Defender 狀態：\n{r.stdout.strip()}"
        elif action == "quick_scan":
            r = subprocess.run(["powershell", "-Command",
                "Start-MpScan -ScanType QuickScan"],
                capture_output=True, text=True, timeout=30)
            return "🛡️ 快速掃描已啟動（背景執行中）"
        elif action == "full_scan":
            r = subprocess.run(["powershell", "-Command",
                "Start-MpScan -ScanType FullScan"],
                capture_output=True, text=True, timeout=30)
            return "🛡️ 完整掃描已啟動（背景執行中）"
        elif action == "threats":
            r = subprocess.run(["powershell", "-Command",
                "Get-MpThreatDetection | Select-Object ThreatID,Resources,ActionSuccess | Format-Table -AutoSize"],
                capture_output=True, text=True, encoding="utf-8", errors="replace")
            out = r.stdout.strip()
            return f"🛡️ 威脅記錄：\n{out}" if out else "✅ 無威脅記錄"
        elif action == "add_exclusion":
            r = subprocess.run(["powershell", "-Command",
                f"Add-MpPreference -ExclusionPath '{path}'"],
                capture_output=True, text=True)
            return f"✅ 已新增排除：{path}" if r.returncode == 0 else f"❌ 失敗：{r.stderr.strip()}"
        elif action == "remove_exclusion":
            r = subprocess.run(["powershell", "-Command",
                f"Remove-MpPreference -ExclusionPath '{path}'"],
                capture_output=True, text=True)
            return f"✅ 已移除排除：{path}" if r.returncode == 0 else f"❌ 失敗：{r.stderr.strip()}"
        elif action == "list_exclusions":
            r = subprocess.run(["powershell", "-Command",
                "(Get-MpPreference).ExclusionPath"],
                capture_output=True, text=True, encoding="utf-8", errors="replace")
            out = r.stdout.strip()
            return f"🛡️ 排除清單：\n{out}" if out else "✅ 無排除項目"
    except subprocess.TimeoutExpired:
        return "⏳ 掃描已啟動（在背景執行）"
    except Exception as e:
        return f"❌ Defender 操作失敗：{e}"


def execute_disk_clean(action="list"):
    try:
        import tempfile, shutil
        tmp = Path(tempfile.gettempdir())
        if action == "list":
            files = list(tmp.rglob("*"))
            total = sum(f.stat().st_size for f in files if f.is_file())
            return f"🗑️ 暫存資料夾：{tmp}\n檔案數：{len(files)}\n佔用：{total/1024/1024:.1f} MB"
        elif action == "clean":
            count = 0
            for f in tmp.iterdir():
                try:
                    if f.is_file(): f.unlink(); count += 1
                    elif f.is_dir(): shutil.rmtree(f, ignore_errors=True); count += 1
                except Exception: pass
            return f"✅ 已清理 {count} 個暫存項目"
    except Exception as e:
        return f"❌ 磁碟清理失敗：{e}"


def execute_drag(x1, y1, x2, y2, dur=0.5):
    pyautogui.moveTo(int(x1), int(y1))
    pyautogui.dragTo(int(x2), int(y2), duration=float(dur), button="left")
    return f"已拖曳 ({x1},{y1}) → ({x2},{y2})"


def execute_firewall(action, name="", port=None, protocol="TCP", direction="Inbound"):
    try:
        if action == "status":
            r = subprocess.run(["powershell", "-Command",
                "Get-NetFirewallProfile | Select-Object Name,Enabled | Format-Table -AutoSize"],
                capture_output=True, text=True, encoding="utf-8", errors="replace")
            return f"🔥 防火牆狀態：\n{r.stdout.strip()}"
        elif action == "list":
            r = subprocess.run(["powershell", "-Command",
                f"Get-NetFirewallRule | Where-Object {{$_.Enabled -eq 'True'}} | Select-Object DisplayName,Direction,Action | Sort-Object Direction | Format-Table -AutoSize"],
                capture_output=True, text=True, encoding="utf-8", errors="replace")
            return f"🔥 防火牆規則：\n{r.stdout.strip()[:2000]}"
        elif action == "add":
            r = subprocess.run(["powershell", "-Command",
                f"New-NetFirewallRule -DisplayName '{name}' -Direction {direction} -Protocol {protocol} -LocalPort {port} -Action Allow -Enabled True"],
                capture_output=True, text=True)
            return f"✅ 已新增防火牆規則：{name} ({direction} {protocol}:{port})" if r.returncode == 0 else f"❌ 失敗：{r.stderr.strip()}"
        elif action == "remove":
            r = subprocess.run(["powershell", "-Command",
                f"Remove-NetFirewallRule -DisplayName '{name}' -Confirm:$false"],
                capture_output=True, text=True)
            return f"✅ 已移除規則：{name}" if r.returncode == 0 else f"❌ 失敗：{r.stderr.strip()}"
        elif action == "enable":
            subprocess.run(["powershell", "-Command",
                "Set-NetFirewallProfile -Profile Domain,Public,Private -Enabled True"], capture_output=True)
            return "✅ 防火牆已啟用"
        elif action == "disable":
            subprocess.run(["powershell", "-Command",
                "Set-NetFirewallProfile -Profile Domain,Public,Private -Enabled False"], capture_output=True)
            return "✅ 防火牆已停用"
    except Exception as e:
        return f"❌ 防火牆操作失敗：{e}"


def execute_font_list(keyword=""):
    try:
        r = subprocess.run(["powershell", "-Command",
            "[System.Reflection.Assembly]::LoadWithPartialName('System.Drawing') | Out-Null; [System.Drawing.FontFamily]::Families | Select-Object -ExpandProperty Name"],
            capture_output=True, text=True, encoding="utf-8", errors="replace")
        fonts = r.stdout.strip().splitlines()
        if keyword:
            fonts = [f for f in fonts if keyword.lower() in f.lower()]
        result = "\n".join(fonts[:50])
        suffix = f"\n...（共 {len(fonts)} 個字型）" if len(fonts) > 50 else f"\n共 {len(fonts)} 個字型"
        return f"🔤 字型清單：\n{result}{suffix}"
    except Exception as e:
        return f"❌ 字型查詢失敗：{e}"


def execute_ftp(action, host="", user="", password="", local="", remote="", port=21):
    try:
        from ftplib import FTP
        ftp = FTP()
        ftp.connect(host, int(port), timeout=30)
        ftp.login(user, password)
        if action == "list":
            items = ftp.nlst(remote or ".")
            ftp.quit()
            return f"📂 FTP {host}{remote or '/'}：\n" + "\n".join(items[:50])
        elif action == "upload":
            with open(local, "rb") as f:
                ftp.storbinary(f"STOR {remote or Path(local).name}", f)
            ftp.quit()
            return f"✅ 已上傳：{local} → {host}/{remote}"
        elif action == "download":
            out = local or str(Path("C:/Users/blue_/Desktop/測試檔案") / Path(remote).name)
            with open(out, "wb") as f:
                ftp.retrbinary(f"RETR {remote}", f.write)
            ftp.quit()
            return f"✅ 已下載：{host}/{remote} → {out}"
        elif action == "delete":
            ftp.delete(remote); ftp.quit()
            return f"✅ 已刪除：{remote}"
        elif action == "mkdir":
            ftp.mkd(remote); ftp.quit()
            return f"✅ 已建立目錄：{remote}"
        elif action == "rename":
            ftp.rename(remote, local); ftp.quit()
            return f"✅ 已重新命名：{remote} → {local}"
    except Exception as e:
        return f"❌ FTP 操作失敗：{e}"


def execute_hotkey(keys: str):
    parts = [k.strip() for k in keys.split("+")]
    pyautogui.hotkey(*parts)
    return f"已執行組合鍵：{keys}"


def execute_hyperv(action, name="", snapshot=""):
    try:
        def ps(cmd):
            r = subprocess.run(["powershell", "-Command", cmd],
                capture_output=True, text=True, encoding="utf-8", errors="replace")
            return r.stdout.strip(), r.returncode
        if action == "list":
            out, _ = ps("Get-VM | Select-Object Name,State,CPUUsage,MemoryAssigned | Format-Table -AutoSize")
            return f"💻 虛擬機清單：\n{out}" if out else "⚠️ 未找到虛擬機（請確認已啟用 Hyper-V）"
        elif action == "status":
            out, _ = ps(f"Get-VM -Name '{name}' | Select-Object * | Format-List")
            return f"💻 {name} 狀態：\n{out}"
        elif action == "start":
            out, rc = ps(f"Start-VM -Name '{name}'")
            return f"✅ 已啟動：{name}" if rc == 0 else f"❌ 啟動失敗：{out}"
        elif action == "stop":
            out, rc = ps(f"Stop-VM -Name '{name}' -Force")
            return f"✅ 已停止：{name}" if rc == 0 else f"❌ 停止失敗：{out}"
        elif action == "pause":
            ps(f"Suspend-VM -Name '{name}'")
            return f"✅ 已暫停：{name}"
        elif action == "resume":
            ps(f"Resume-VM -Name '{name}'")
            return f"✅ 已繼續：{name}"
        elif action == "snapshot":
            sname = snapshot or dt.dt.datetime.now().strftime("snap_%Y%m%d_%H%M%S")
            out, rc = ps(f"Checkpoint-VM -Name '{name}' -SnapshotName '{sname}'")
            return f"✅ 快照已建立：{sname}" if rc == 0 else f"❌ 失敗：{out}"
        elif action == "restore":
            out, rc = ps(f"Restore-VMSnapshot -Name '{name}' -VMName '{name}' -Confirm:$false")
            return f"✅ 已還原快照：{snapshot}" if rc == 0 else f"❌ 失敗：{out}"
        elif action == "delete_snapshot":
            ps(f"Remove-VMSnapshot -VMName '{name}' -Name '{snapshot}' -Confirm:$false")
            return f"✅ 已刪除快照：{snapshot}"
    except Exception as e:
        return f"❌ Hyper-V 操作失敗：{e}"


def execute_interval_schedule(action, name="", command="", every_minutes=60.0, repeat=0, duration_hours=0.0):
    global _interval_schedules
    try:
        if action == "list":
            if not _interval_schedules:
                return "⚠️ 無執行中排程"
            return "⏱️ 間隔排程：\n" + "\n".join(f"- {k}: 每{v['mins']}分鐘，已執行{v['count']}次" for k,v in _interval_schedules.items())
        elif action == "stop":
            if name in _interval_schedules:
                _interval_schedules[name]["running"] = False
                del _interval_schedules[name]
                return f"✅ 已停止排程：{name}"
            return f"⚠️ 找不到排程：{name}"
        elif action == "start":
            if name in _interval_schedules:
                return f"⚠️ 已有同名排程：{name}"
            cfg = {"command": command, "mins": every_minutes, "repeat": repeat, "count": 0, "running": True}
            _interval_schedules[name] = cfg
            def _run():
                import time as t
                end_time = t.time() + float(duration_hours) * 3600 if duration_hours > 0 else float("inf")
                max_count = int(repeat) if repeat > 0 else float("inf")
                while _interval_schedules.get(name, {}).get("running"):
                    if t.time() > end_time or _interval_schedules.get(name, {}).get("count", 0) >= max_count:
                        _interval_schedules.pop(name, None); break
                    subprocess.Popen(command, shell=True)
                    _interval_schedules[name]["count"] = _interval_schedules[name].get("count", 0) + 1
                    t.sleep(float(every_minutes) * 60)
            threading.Thread(target=_run, daemon=True).start()
            desc = f"每 {every_minutes} 分鐘" + (f"，共 {repeat} 次" if repeat else "") + (f"，持續 {duration_hours} 小時" if duration_hours else "")
            return f"✅ 間隔排程已啟動：{name}（{desc}）"
    except Exception as e:
        return f"❌ 間隔排程失敗：{e}"


def execute_lock_screen(action):
    try:
        if action == "lock":
            subprocess.Popen(["rundll32.exe", "user32.dll,LockWorkStation"])
            return "🔒 螢幕已鎖定"
        elif action == "logoff":
            subprocess.run(["shutdown", "/l"], capture_output=True)
            return "✅ 已登出"
        elif action == "switch_user":
            subprocess.Popen(["tsdiscon.exe"])
            return "✅ 已切換使用者"
    except Exception as e:
        return f"❌ 鎖定/登出失敗：{e}"


def execute_net_share(action, share_path="", drive="Z:", user="", password=""):
    try:
        if action == "list":
            r = subprocess.run(["net", "use"], capture_output=True, text=True, encoding="cp950", errors="replace")
            return f"🌐 網路磁碟機：\n{r.stdout.strip()}"
        elif action == "connect":
            args = ["net", "use", drive, share_path]
            if user: args += [f"/user:{user}", password]
            r = subprocess.run(args, capture_output=True, text=True, encoding="cp950", errors="replace")
            return r.stdout.strip() or f"✅ 已連線 {share_path} → {drive}"
        elif action == "disconnect":
            r = subprocess.run(["net", "use", drive, "/delete"], capture_output=True, text=True, encoding="cp950", errors="replace")
            return r.stdout.strip() or f"✅ 已中斷 {drive}"
    except Exception as e:
        return f"❌ 網路芳鄰失敗：{e}"


def execute_object_detect(target, action="find", region=""):
    try:
        import anthropic, base64, io as _io, json, re
        reg = None
        if region:
            parts = [int(v) for v in region.split(",")]
            if len(parts) == 4: reg = tuple(parts)
        screenshot = pyautogui.screenshot(region=reg)
        buf = _io.BytesIO(); screenshot.save(buf, format="PNG")
        img_b64 = base64.standard_b64encode(buf.getvalue()).decode()
        offset_x = reg[0] if reg else 0
        offset_y = reg[1] if reg else 0
        _client = anthropic.Anthropic(api_key=os.getenv("ANTHROPIC_API_KEY"))
        resp = _client.messages.create(
            model="claude-sonnet-4-6", max_tokens=256,
            messages=[{"role": "user", "content": [
                {"type": "image", "source": {"type": "base64", "media_type": "image/png", "data": img_b64}},
                {"type": "text", "text": f"在圖片中找到「{target}」的位置。回答 JSON：{{\"found\": true/false, \"x\": 中心X座標, \"y\": 中心Y座標, \"description\": \"說明\"}}"}
            ]}]
        )
        m = re.search(r'\{.*\}', resp.content[0].text, re.DOTALL)
        if not m: return f"⚠️ AI 無法解析回應"
        result = json.loads(m.group())
        if not result.get("found"): return f"⚠️ 未找到：{target}"
        ax = result["x"] + offset_x
        ay = result["y"] + offset_y
        if action == "click": pyautogui.click(ax, ay); return f"✅ 已點擊「{target}」({ax},{ay})"
        elif action == "double_click": pyautogui.doubleClick(ax, ay); return f"✅ 已雙擊「{target}」({ax},{ay})"
        return f"✅ 找到「{target}」：座標({ax},{ay}) — {result.get('description','')}"
    except Exception as e:
        return f"❌ 物件偵測失敗：{e}"


def execute_onedrive(action, path="", remote=""):
    try:
        import shutil
        onedrive_path = os.path.expandvars(r"%USERPROFILE%\OneDrive")
        if not Path(onedrive_path).exists():
            onedrive_path = os.path.expandvars(r"%OneDrive%") or str(Path.home() / "OneDrive")
        if action == "status":
            r = subprocess.run(["powershell", "-Command",
                "Get-Process OneDrive -ErrorAction SilentlyContinue | Select-Object Name,CPU,WorkingSet"],
                capture_output=True, text=True, encoding="utf-8", errors="replace")
            size = sum(f.stat().st_size for f in Path(onedrive_path).rglob("*") if f.is_file()) / 1024 / 1024 / 1024
            return f"☁️ OneDrive 路徑：{onedrive_path}\n使用空間：{size:.2f} GB\n程序狀態：\n{r.stdout.strip()}"
        elif action == "list":
            target = Path(onedrive_path) / (remote or "")
            if not target.exists(): return f"⚠️ 路徑不存在：{target}"
            items = list(target.iterdir())
            lines = [f"{'📁' if p.is_dir() else '📄'} {p.name}" for p in sorted(items)]
            return f"☁️ OneDrive/{remote or ''}：\n" + "\n".join(lines[:30])
        elif action == "upload":
            dest = Path(onedrive_path) / (remote or Path(path).name)
            shutil.copy2(path, dest)
            return f"✅ 已上傳至 OneDrive：{dest}"
        elif action == "download":
            src = Path(onedrive_path) / remote
            if not src.exists(): return f"⚠️ 找不到：{src}"
            dest = path or str(Path("C:/Users/blue_/Desktop/測試檔案") / src.name)
            shutil.copy2(src, dest)
            return f"✅ 已從 OneDrive 下載：{dest}"
        elif action == "sync":
            subprocess.Popen(["powershell", "-Command",
                'Start-Process "$env:LOCALAPPDATA\\Microsoft\\OneDrive\\OneDrive.exe"'])
            return "✅ OneDrive 同步已觸發"
        elif action == "open":
            os.startfile(onedrive_path)
            return f"✅ 已開啟 OneDrive 資料夾：{onedrive_path}"
    except Exception as e:
        return f"❌ OneDrive 操作失敗：{e}"


def execute_proxy(action, host=""):
    try:
        import winreg
        key_path = r"Software\Microsoft\Windows\CurrentVersion\Internet Settings"
        with winreg.OpenKey(winreg.HKEY_CURRENT_USER, key_path, 0, winreg.KEY_READ | winreg.KEY_SET_VALUE) as k:
            if action == "get":
                try:
                    enabled = winreg.QueryValueEx(k, "ProxyEnable")[0]
                    server = winreg.QueryValueEx(k, "ProxyServer")[0]
                    return f"🌐 代理設定：{'啟用' if enabled else '停用'}\n伺服器：{server}"
                except Exception:
                    return "🌐 代理：未設定"
            elif action == "set":
                winreg.SetValueEx(k, "ProxyEnable", 0, winreg.REG_DWORD, 1)
                winreg.SetValueEx(k, "ProxyServer", 0, winreg.REG_SZ, host)
                return f"✅ 代理已設定：{host}"
            elif action == "disable":
                winreg.SetValueEx(k, "ProxyEnable", 0, winreg.REG_DWORD, 0)
                return "✅ 代理已停用"
    except Exception as e:
        return f"❌ 代理設定失敗：{e}"


def execute_rdp_connect(host, user="", width=1280, height=720):
    try:
        args = ["/v:" + host, f"/w:{width}", f"/h:{height}"]
        if user:
            args.append(f"/u:{user}")
        subprocess.Popen(["mstsc"] + args)
        return f"✅ 正在連線 RDP：{host}"
    except Exception as e:
        return f"❌ RDP 連線失敗：{e}"


def execute_right_menu(x, y, item=""):
    try:
        pyautogui.rightClick(int(x), int(y))
        time.sleep(0.3)
        if item:
            pyautogui.write(item, interval=0.05)
            time.sleep(0.2)
            pyautogui.press("enter")
            return f"✅ 已右鍵點擊並選擇：{item}"
        return f"✅ 已右鍵點擊 ({x},{y})"
    except Exception as e:
        return f"❌ 右鍵選單失敗：{e}"


def execute_screen_live(action, fps=0.5, duration=60.0, quality=50, _bot_send=None, _chat_id=None):
    global _screen_live_running
    try:
        if action == "stop":
            _screen_live_running = False
            return "✅ 螢幕串流已停止"
        elif action == "start":
            if _screen_live_running:
                return "⚠️ 螢幕串流已在執行中"
            _screen_live_running = True
            def _stream():
                global _screen_live_running
                import io as _io, time as t, asyncio
                interval = 1.0 / max(float(fps), 0.1)
                end = t.time() + float(duration)
                try:
                    loop = asyncio.get_event_loop()
                except Exception:
                    loop = None
                count = 0
                while _screen_live_running and t.time() < end:
                    try:
                        screenshot = pyautogui.screenshot()
                        buf = _io.BytesIO()
                        screenshot.save(buf, format="JPEG", quality=int(quality))
                        buf.seek(0)
                        if _bot_send and _chat_id and loop:
                            import telegram
                            asyncio.run_coroutine_threadsafe(
                                _bot_send(_chat_id, photo=buf),
                                loop)
                        count += 1
                    except Exception:
                        pass
                    t.sleep(interval)
                _screen_live_running = False
            threading.Thread(target=_stream, daemon=True).start()
            return f"✅ 螢幕串流已啟動（{fps} FPS，{duration}s，畫質 {quality}）"
    except Exception as e:
        return f"❌ 螢幕串流失敗：{e}"


def execute_screen_stream(duration=10, interval=2):
    screenshots = []
    end = time.time() + duration
    while time.time() < end:
        img = pyautogui.screenshot()
        buf = io.BytesIO()
        img.save(buf, format="PNG")
        screenshots.append(buf.getvalue())
        time.sleep(interval)
    return screenshots


def execute_screen_watch(template_path, command, timeout=60):
    import time as t
    start = t.time()
    while t.time() - start < timeout:
        try:
            loc = pyautogui.locateOnScreen(template_path, confidence=0.8)
            if loc:
                subprocess.run(command, shell=True)
                return f"偵測到目標，已執行：{command}"
        except Exception:
            pass
        t.sleep(2)
    return "監控逾時，未偵測到目標"


def execute_sysres_chart(duration=10):
    try:
        import psutil, matplotlib
        matplotlib.use("Agg")
        import matplotlib.pyplot as plt, time as t
        cpu_vals, mem_vals = [], []
        for _ in range(duration):
            cpu_vals.append(psutil.cpu_percent(interval=1))
            mem_vals.append(psutil.virtual_memory().percent)
        out = str(Path("C:/Users/blue_/Desktop/測試檔案") / f"sysres_{dt.dt.datetime.now().strftime('%H%M%S')}.png")
        fig, ax = plt.subplots()
        ax.plot(range(1,duration+1), cpu_vals, label="CPU %", color="blue")
        ax.plot(range(1,duration+1), mem_vals, label="RAM %", color="orange")
        ax.set_ylim(0,100); ax.legend(); ax.set_title("系統資源使用率")
        plt.tight_layout(); plt.savefig(out); plt.close()
        return out
    except Exception as e:
        return f"❌ 失敗：{e}"


def execute_usb_list():
    try:
        r = subprocess.run(["powershell", "-Command",
            "Get-PnpDevice | Where-Object {$_.Class -eq 'USB' -and $_.Status -eq 'OK'} | Select-Object FriendlyName,InstanceId | Format-Table -AutoSize"],
            capture_output=True, text=True, encoding="utf-8", errors="replace")
        return f"🔌 USB 裝置：\n{r.stdout.strip()[:2000]}" if r.stdout.strip() else "⚠️ 無 USB 裝置"
    except Exception as e:
        return f"❌ USB 查詢失敗：{e}"


def execute_wake_on_lan(mac, broadcast="255.255.255.255", port=9):
    try:
        import socket
        mac_clean = mac.replace(":", "").replace("-", "")
        if len(mac_clean) != 12:
            return f"❌ MAC 位址格式錯誤：{mac}"
        mac_bytes = bytes.fromhex(mac_clean)
        magic = b"\xff" * 6 + mac_bytes * 16
        with socket.socket(socket.AF_INET, socket.SOCK_DGRAM) as sock:
            sock.setsockopt(socket.SOL_SOCKET, socket.SO_BROADCAST, 1)
            sock.sendto(magic, (broadcast, int(port)))
        return f"✅ WOL 魔法封包已送出 → {mac}（{broadcast}:{port}）"
    except Exception as e:
        return f"❌ WOL 失敗：{e}"


def execute_wifi(action, ssid="", password=""):
    try:
        if action == "scan":
            r = subprocess.run(["netsh", "wlan", "show", "networks", "mode=Bssid"],
                capture_output=True, text=True, encoding="utf-8", errors="replace")
            return f"📡 附近 Wi-Fi：\n{r.stdout.strip()[:2000]}"
        elif action == "status":
            r = subprocess.run(["netsh", "wlan", "show", "interfaces"],
                capture_output=True, text=True, encoding="utf-8", errors="replace")
            return f"📡 Wi-Fi 狀態：\n{r.stdout.strip()}"
        elif action == "saved":
            r = subprocess.run(["netsh", "wlan", "show", "profiles"],
                capture_output=True, text=True, encoding="utf-8", errors="replace")
            return f"📡 已儲存 Wi-Fi：\n{r.stdout.strip()}"
        elif action == "password":
            r = subprocess.run(["netsh", "wlan", "show", "profile", f"name={ssid}", "key=clear"],
                capture_output=True, text=True, encoding="utf-8", errors="replace")
            return f"🔑 '{ssid}' 密碼資訊：\n{r.stdout.strip()}"
        elif action == "connect":
            r = subprocess.run(["netsh", "wlan", "connect", f"name={ssid}"],
                capture_output=True, text=True, encoding="utf-8", errors="replace")
            return f"✅ 嘗試連線 '{ssid}'：{r.stdout.strip()}"
        elif action == "disconnect":
            r = subprocess.run(["netsh", "wlan", "disconnect"],
                capture_output=True, text=True, encoding="utf-8", errors="replace")
            return f"✅ Wi-Fi 已斷線：{r.stdout.strip()}"
    except Exception as e:
        return f"❌ Wi-Fi 操作失敗：{e}"


def execute_wifi_hotspot(action, ssid="", password=""):
    try:
        if action == "status":
            r = subprocess.run(["powershell", "-Command",
                "Get-NetConnectionProfile | Select-Object Name,NetworkCategory | Format-Table; netsh wlan show hostednetwork"],
                capture_output=True, text=True, encoding="utf-8", errors="replace")
            return f"📡 熱點狀態：\n{r.stdout.strip()}"
        elif action == "set":
            r = subprocess.run(["netsh", "wlan", "set", "hostednetwork",
                f"mode=allow", f"ssid={ssid}", f"key={password}"],
                capture_output=True, text=True, encoding="utf-8", errors="replace")
            return f"✅ 熱點設定完成：SSID={ssid}\n{r.stdout.strip()}"
        elif action == "start":
            r = subprocess.run(["netsh", "wlan", "start", "hostednetwork"],
                capture_output=True, text=True, encoding="utf-8", errors="replace")
            return f"✅ 熱點已啟動\n{r.stdout.strip()}"
        elif action == "stop":
            r = subprocess.run(["netsh", "wlan", "stop", "hostednetwork"],
                capture_output=True, text=True, encoding="utf-8", errors="replace")
            return f"✅ 熱點已停止\n{r.stdout.strip()}"
    except Exception as e:
        return f"❌ WiFi 熱點失敗：{e}"


def execute_win_service(action, name=""):
    try:
        if action == "list":
            r = subprocess.run(["powershell.exe", "-Command",
                "Get-Service | Select-Object Name,Status | Format-Table -AutoSize"],
                capture_output=True, text=True, encoding="utf-8", errors="replace")
            return r.stdout[:2000]
        else:
            cmd = f"{'Start' if action=='start' else 'Stop'}-Service -Name '{name}' -Force"
            r = subprocess.run(["powershell.exe", "-Command", cmd],
                capture_output=True, text=True, encoding="utf-8", errors="replace")
            return r.stdout or r.stderr or f"✅ {action} {name}"
    except Exception as e:
        return f"❌ 服務操作失敗：{e}"


def execute_wsl(action, distro="", command=""):
    try:
        if action == "list":
            r = subprocess.run(["wsl", "--list", "--verbose"],
                capture_output=True, text=True, encoding="utf-16-le", errors="replace")
            return f"🐧 WSL 發行版：\n{r.stdout.strip()}"
        elif action == "status":
            r = subprocess.run(["wsl", "--status"],
                capture_output=True, text=True, encoding="utf-16-le", errors="replace")
            return f"🐧 WSL 狀態：\n{r.stdout.strip()}"
        elif action == "run":
            cmd = ["wsl"]
            if distro: cmd += ["-d", distro]
            cmd += ["--", "bash", "-c", command]
            r = subprocess.run(cmd, capture_output=True, text=True, encoding="utf-8", errors="replace")
            return f"🐧 WSL 輸出：\n{r.stdout.strip()}" + (f"\n錯誤：{r.stderr.strip()}" if r.stderr.strip() else "")
        elif action == "start":
            subprocess.Popen(["wsl"] + (["-d", distro] if distro else []))
            return f"✅ WSL 已啟動：{distro or '預設'}"
        elif action == "stop":
            cmd = ["wsl", "--terminate", distro] if distro else ["wsl", "--shutdown"]
            subprocess.run(cmd, capture_output=True)
            return f"✅ WSL 已停止：{distro or '全部'}"
        elif action == "install":
            r = subprocess.run(["wsl", "--install", "-d", distro],
                capture_output=True, text=True)
            return f"✅ 正在安裝 {distro}（需要網路）"
    except Exception as e:
        return f"❌ WSL 操作失敗：{e}（請確認已啟用 WSL）"


# ── 主程式 ──────────────────────────────────────────

if __name__ == "__main__":
    if len(sys.argv) < 2:
        print(__doc__)
        sys.exit(0)

    tool = sys.argv[1]
    args = sys.argv[2:]

    tools = {
        "weather":       lambda: get_weather(args[0]),
        "stock":         lambda: get_stock(args[0], args[1] if len(args) > 1 else "1mo"),
        "image":         lambda: generate_image(args[0], args[1] if len(args) > 1 else ""),
        "screenshot":    lambda: screenshot(),
        "click":         lambda: click(*args),
        "double_click":  lambda: double_click(*args),
        "right_click":   lambda: right_click(*args),
        "move":          lambda: move(*args),
        "type":          lambda: type_text(" ".join(args)),
        "press":         lambda: press_key(args[0]),
        "open":          lambda: open_app(" ".join(args)),
        "scroll":        lambda: scroll(*args),
        "pos":           lambda: pos(),
        "bot_status":    lambda: bot_status(),
        "bot_restart":   lambda: bot_restart(),
        "schedule_list": lambda: schedule_list(),
        "schedule_add":  lambda: schedule_add(args[0], args[1], args[2]),
        "schedule_del":  lambda: schedule_del(args[0]),
        "memory_save":   lambda: memory_save(int(args[0]), " ".join(args[1:])),
        "memory_list":   lambda: memory_list(int(args[0])) if args else print("用法：memory_list <chat_id>"),
        "memory_del":    lambda: memory_del(int(args[0]), int(args[1])),
        "vision":        lambda: vision(" ".join(args) if args else "請描述這個畫面上有什麼，以及目前電腦在做什麼事。"),
        "find_image":    lambda: find_image(args[0], float(args[1]) if len(args) > 1 else 0.8),
        "browser":       lambda: browser(args[0], *args[1:]),
        "window_list":   lambda: window_list(),
        "window_focus":  lambda: window_focus(" ".join(args)),
        "window_close":  lambda: window_close(" ".join(args)),
        "window_min":    lambda: window_min(" ".join(args)),
        "window_max":    lambda: window_max(" ".join(args)),
        "hotkey":        lambda: hotkey(*args),
        "clipboard_get": lambda: clipboard_get(),
        "clipboard_set": lambda: clipboard_set(" ".join(args)),
        "file_list":     lambda: file_list(args[0] if args else "."),
        "file_read":     lambda: file_read(args[0]),
        "file_write":    lambda: file_write(args[0], " ".join(args[1:])),
        "file_delete":   lambda: file_delete(args[0]),
        "file_copy":     lambda: file_copy(args[0], args[1]),
        "file_move":     lambda: file_move(args[0], args[1]),
        "file_search":   lambda: file_search(args[0], args[1]),
        "sysinfo":       lambda: sysinfo(),
        "process_list":  lambda: process_list(),
        "process_kill":  lambda: process_kill(args[0]),
        "notify":        lambda: notify(args[0], " ".join(args[1:])),
        "tts":           lambda: tts(" ".join(args)),
        "record_start":  lambda: record_start(),
        "record_stop":   lambda: record_stop(),
        "record_play":   lambda: record_play(args[0]),
        "email":         lambda: send_email(args[0], args[1], " ".join(args[2:])),
        "stt":           lambda: stt(args[0] if args and args[0].isdigit() else 5, args[0] if args and not args[0].isdigit() else "", args[1] if len(args)>1 else "zh-TW"),
        "ocr":           lambda: ocr(args[0] if args else ""),
        "workflow_run":  lambda: workflow_run(args[0]),
        "workflow_save": lambda: workflow_save(args[0], " ".join(args[1:])),
        "screen_watch":  lambda: screen_watch(args[0], args[1], float(args[2]) if len(args)>2 else 2.0),
        "monitors":      lambda: monitors(),
        "zip":           lambda: zip_files(args[0], args[1]),
        "unzip":         lambda: unzip(args[0], args[1]),
        "download":      lambda: download(args[0], args[1] if len(args)>1 else ""),
        "print_file":    lambda: print_file(args[0]),
        "wifi_list":       lambda: wifi_list(),
        "wifi_connect":    lambda: wifi_connect(args[0], args[1]),
        "screen_stream":   lambda: screen_stream(int(args[0]) if args else 10),
        "wake_listen":     lambda: wake_listen(args[0] if args else "小牛馬"),
        "drag":            lambda: drag(*args),
        "right_menu":      lambda: right_menu(args[0], args[1], args[2] if len(args)>2 else ""),
        "ai_plan":         lambda: ai_plan(" ".join(args)),
        "clipboard_history": lambda: clipboard_history(),
        "vdesktop":        lambda: vdesktop(args[0]),
        "power":           lambda: power(args[0]),
        "bt_scan":         lambda: bt_scan(),
        "bt_connect":      lambda: bt_connect(args[0]),
        "run_python":      lambda: run_python(" ".join(args)),
        "run_shell":       lambda: run_shell(" ".join(args)),
        "word_read":       lambda: word_read(args[0]),
        "word_write":      lambda: word_write(args[0], " ".join(args[1:])),
        "excel_read":      lambda: excel_read(args[0], args[1] if len(args)>1 else None),
        "excel_write":     lambda: excel_write(args[0], args[1], " ".join(args[2:])),
        "pdf_read":        lambda: pdf_read(args[0]),
        "screen_diff":     lambda: screen_diff(float(args[0]) if args else 2.0, args[1] if len(args)>1 else "full"),
        "scrape":          lambda: scrape(args[0], args[1] if len(args)>1 else "body"),
        "img_edit":        lambda: img_edit(args[0], args[1], *args[2:]),
        "gdrive_upload":   lambda: gdrive_upload(args[0], args[1] if len(args)>1 else "root"),
        "gdrive_download": lambda: gdrive_download(args[0], args[1]),
        "db_query":        lambda: db_query(args[0], " ".join(args[1:])),
        "db_mysql":        lambda: db_mysql(args[0], args[1], " ".join(args[2:])),
        "encrypt":         lambda: encrypt(args[0], args[1]),
        "decrypt":         lambda: decrypt(args[0], args[1]),
        "clipboard_watch": lambda: clipboard_watch(float(args[0]) if args else 30.0),
        "qr_gen":          lambda: qr_gen(args[0], args[1] if len(args)>1 else ""),
        "qr_scan":         lambda: qr_scan(args[0] if args else ""),
        "screen_record":   lambda: screen_record(float(args[0]) if args else 10.0, args[1] if len(args)>1 else ""),
        "webcam":          lambda: webcam_capture(args[0] if args else ""),
        "translate":       lambda: translate(" ".join(args[:1] if len(args)==1 else [args[0]]), args[1] if len(args)>1 else "zh-TW", args[2] if len(args)>2 else "auto"),
        "chart":           lambda: chart(args[0], args[1], args[2] if len(args)>2 else "", args[3] if len(args)>3 else ""),
        "pptx_read":       lambda: pptx_read(args[0]),
        "pptx_create":     lambda: pptx_create(args[0], " ".join(args[1:])),
        "api_call":        lambda: api_call(args[0], args[1], args[2] if len(args)>2 else "{}", args[3] if len(args)>3 else "{}"),
        "watchdog":        lambda: watchdog(args[0], args[1], float(args[2]) if len(args)>2 else 60.0),
        "ssh_run":         lambda: ssh_run(args[0], args[1], args[2], " ".join(args[3:])),
        "sftp_upload":     lambda: sftp_upload(args[0], args[1], args[2], args[3], args[4]),
        "sftp_download":   lambda: sftp_download(args[0], args[1], args[2], args[3], args[4]),
        "net_ping":        lambda: net_ping(args[0], int(args[1]) if len(args)>1 else 4),
        "net_traceroute":  lambda: net_traceroute(args[0]),
        "net_portscan":    lambda: net_portscan(args[0], args[1] if len(args)>1 else "22,80,443,3306,3389,8080"),
        "win_service":     lambda: win_service(args[0], args[1] if len(args)>1 else ""),
        "pdf_merge":       lambda: pdf_merge(args[0], args[1]),
        "pdf_split":       lambda: pdf_split(args[0], args[1]),
        "pdf_watermark":   lambda: pdf_watermark(args[0], args[1], args[2] if len(args)>2 else ""),
        "audio_convert":   lambda: audio_convert(args[0], args[1]),
        "audio_trim":      lambda: audio_trim(args[0], int(args[1]), int(args[2]), args[3] if len(args)>3 else ""),
        "discord_notify":  lambda: discord_notify(args[0], " ".join(args[1:])),
        "line_notify":     lambda: line_notify(args[0], " ".join(args[1:])),
        "disk_clean":      lambda: disk_clean(args[0] if args else "list"),
        "backup":          lambda: backup(args[0], args[1]),
        "registry_read":   lambda: registry_read(args[0], args[1] if len(args)>1 else ""),
        "registry_write":  lambda: registry_write(args[0], args[1], " ".join(args[2:])),
        "video_screenshot":lambda: video_screenshot(args[0], float(args[1]) if len(args)>1 else 0, args[2] if len(args)>2 else ""),
        "video_trim":      lambda: video_trim(args[0], float(args[1]), float(args[2]), args[3] if len(args)>3 else ""),
        "monitor_list":    lambda: monitor_list(),
        "email_read":      lambda: email_read(args[0], args[1], args[2], args[3] if len(args)>3 else "INBOX", int(args[4]) if len(args)>4 else 5),
        "gcal_list":       lambda: gcal_list(int(args[0]) if args else 7),
        "gcal_add":        lambda: gcal_add(args[0], args[1], args[2], args[3] if len(args)>3 else ""),
        "global_hotkey":   lambda: global_hotkey_listen(args[0], args[1], float(args[2]) if len(args)>2 else 60.0),
        "git_op":          lambda: git_op(args[0], args[1] if len(args)>1 else ".", args[2] if len(args)>2 else "", args[3] if len(args)>3 else "origin", args[4] if len(args)>4 else "master"),
        "hw_monitor":      lambda: hw_monitor(),
        "report_gen":      lambda: report_gen(args[0], args[1], args[2] if len(args)>2 else ""),
        "dropbox_upload":  lambda: dropbox_upload(args[0], args[1], args[2] if len(args)>2 else ""),
        "dropbox_download":lambda: dropbox_download(args[0], args[1], args[2] if len(args)>2 else ""),
        "docker_op":       lambda: docker_op(args[0], args[1] if len(args)>1 else ""),
        "pdf_to_images":   lambda: pdf_to_images(args[0], args[1] if len(args)>1 else "", int(args[2]) if len(args)>2 else 150),
        "barcode_scan":    lambda: barcode_scan(args[0] if args else ""),
        "nlp_summarize":   lambda: nlp_summarize(" ".join(args)),
        "nlp_sentiment":   lambda: nlp_sentiment(" ".join(args)),
        "vpn_control":     lambda: vpn_control(args[0], args[1] if len(args)>1 else "", args[2] if len(args)>2 else "", args[3] if len(args)>3 else ""),
        "sys_restore":     lambda: sys_restore(args[0], " ".join(args[1:]) if len(args)>1 else ""),
        "disk_analyze":    lambda: disk_analyze(args[0] if args else "C:/", int(args[1]) if len(args)>1 else 10),
        "face_detect":     lambda: face_detect(args[0] if args else "", args[1] if len(args)>1 else ""),
        "video_to_gif":    lambda: video_to_gif(args[0], float(args[1]) if len(args)>1 else 0, float(args[2]) if len(args)>2 else 5.0, args[3] if len(args)>3 else "", int(args[4]) if len(args)>4 else 10),
        "excel_chart":     lambda: excel_chart(args[0], args[1], args[2] if len(args)>2 else "bar", args[3] if len(args)>3 else ""),
        "speedtest":       lambda: speedtest_run(),
        "screenshot_compare": lambda: screenshot_compare(args[0] if args else "", args[1] if len(args)>1 else "", args[2] if len(args)>2 else ""),
        "set_reminder":    lambda: set_reminder(args[0], " ".join(args[1:])),
        "webpage_screenshot": lambda: webpage_screenshot(args[0], args[1] if len(args)>1 else ""),
        "web_monitor":     lambda: web_monitor(args[0], args[1] if len(args)>1 else "body", float(args[2]) if len(args)>2 else 60.0, float(args[3]) if len(args)>3 else 3600.0),
        "batch_rename":    lambda: batch_rename(args[0], args[1], args[2], args[3] if len(args)>3 else ""),
        "img_compress":    lambda: img_compress(args[0], int(args[1]) if len(args)>1 else 75, args[2] if len(args)>2 else ""),
        "batch_img_process": lambda: batch_img_process(args[0], args[1], int(args[2]) if len(args)>2 else 0, int(args[3]) if len(args)>3 else 0, int(args[4]) if len(args)>4 else 75),
        "ocr_translate":   lambda: ocr_translate(args[0] if args else "", args[1] if len(args)>1 else "zh-TW"),
        "ip_info":         lambda: ip_info(args[0] if args else ""),
        "currency":        lambda: currency(float(args[0]), args[1], args[2]),
        "event_log":       lambda: event_log(args[0] if args else "System", args[1] if len(args)>1 else "Error", int(args[2]) if len(args)>2 else 10),
        "tts_edge":        lambda: tts_edge(" ".join(args[:-1]) if len(args)>1 else " ".join(args), args[-1] if len(args)>1 and args[-1].startswith("zh-") else "zh-TW-HsiaoChenNeural"),
        "tts_voices":      lambda: tts_voices(),
        "send_email_attach": lambda: send_email_attach(args[0], args[1], args[2], args[3] if len(args)>3 else ""),
        "clipboard_img_get": lambda: clipboard_img_get(args[0] if args else ""),
        "clipboard_img_set": lambda: clipboard_img_set(args[0]),
        "usb_list":        lambda: usb_list(),
        "firewall":        lambda: firewall(args[0], args[1] if len(args)>1 else "", args[2] if len(args)>2 else "in", int(args[3]) if len(args)>3 else 0, args[4] if len(args)>4 else "TCP"),
        "todo":            lambda: todo(args[0], " ".join(args[1:]) if args[0]=="add" else "", int(args[1]) if len(args)>1 and args[0] in ("done","delete") else 0),
        "file_sync":       lambda: file_sync(args[0], args[1], args[2].lower()=="true" if len(args)>2 else False),
        "sysres_chart":    lambda: sysres_chart(int(args[0]) if args else 10, args[1] if len(args)>1 else ""),
        "password_save":   lambda: password_save(args[0], args[1], args[2], args[3]),
        "password_get":    lambda: password_get(args[0], args[1]),
        "rdp_connect":     lambda: rdp_connect(args[0], args[1] if len(args)>1 else "", int(args[2]) if len(args)>2 else 1280, int(args[3]) if len(args)>3 else 720),
        "chrome_bookmarks": lambda: chrome_bookmarks(),
        "printer_list":    lambda: printer_list(),
        "printer_jobs":    lambda: printer_jobs(),
        "net_share":       lambda: net_share(args[0], args[1] if len(args)>1 else "", args[2] if len(args)>2 else "Z:", args[3] if len(args)>3 else "", args[4] if len(args)>4 else ""),
        "font_list":       lambda: font_list(args[0] if args else ""),
        "volume_get":      lambda: volume_get(),
        "volume_set":      lambda: volume_set(int(args[0])),
        "volume_mute":     lambda: volume_mute(args[0].lower() != "false" if args else True),
        "brightness_get":  lambda: brightness_get(),
        "brightness_set":  lambda: brightness_set(int(args[0])),
        "resolution_list": lambda: resolution_list(),
        "media_control":   lambda: media_control(args[0]),
        "audio_devices":   lambda: audio_devices(),
        "audio_switch":    lambda: audio_switch(" ".join(args)),
        "software_list":   lambda: software_list(args[0] if args else ""),
        "software_install":lambda: software_install(" ".join(args)),
        "software_uninstall": lambda: software_uninstall(" ".join(args)),
        "startup_list":    lambda: startup_list(),
        "startup_add":     lambda: startup_add(args[0], " ".join(args[1:])),
        "startup_remove":  lambda: startup_remove(args[0]),
        "env_get":         lambda: env_get(args[0] if args else ""),
        "env_set":         lambda: env_set(args[0], args[1], len(args)>2 and args[2].lower()=="true"),
        "user_list":       lambda: user_list(),
        "user_create":     lambda: user_create(args[0], args[1]),
        "user_delete":     lambda: user_delete(args[0]),
        "win_update":      lambda: win_update(args[0] if args else "list"),
        "device_list":     lambda: device_list(args[0] if args else ""),
        "device_toggle":   lambda: device_toggle(args[0], len(args)<2 or args[1].lower()!="false"),
        "netadapter_list": lambda: netadapter_list(),
        "netadapter_toggle": lambda: netadapter_toggle(args[0], len(args)<2 or args[1].lower()!="false"),
        "dns_get":         lambda: dns_get(),
        "dns_set":         lambda: dns_set(args[0], args[1], args[2] if len(args)>2 else ""),
        "ip_config":       lambda: ip_config(args[0], args[1], args[2] if len(args)>2 else "255.255.255.0", args[3] if len(args)>3 else ""),
        "hosts_list":      lambda: hosts_list(),
        "hosts_add":       lambda: hosts_add(args[0], args[1]),
        "hosts_remove":    lambda: hosts_remove(args[0]),
        "net_traffic":     lambda: net_traffic(int(args[0]) if args else 5),
        "if_then":         lambda: if_then(args[0], args[1], " ".join(args[2:-1]) if len(args)>3 else args[2], float(args[-1]) if len(args)>3 else 300.0),
        "window_arrange":  lambda: window_arrange(args[0] if args else "side_by_side"),
        "region_ocr":      lambda: region_ocr(int(args[0]), int(args[1]), int(args[2]), int(args[3]), args[4] if len(args)>4 else "ch_tra"),
        "window_screenshot": lambda: window_screenshot(args[0], args[1] if len(args)>1 else ""),
        "firewall":          lambda: firewall(args[0], args[1] if len(args)>1 else "", int(args[2]) if len(args)>2 else None, args[3] if len(args)>3 else "TCP", args[4] if len(args)>4 else "Inbound"),
        "process_mgr":       lambda: process_mgr(args[0], args[1] if len(args)>1 else "", int(args[2]) if len(args)>2 and args[2].isdigit() else None, args[3] if len(args)>3 else "normal"),
        "power_plan":        lambda: power_plan(args[0], args[1] if len(args)>1 else "balanced"),
        "event_log":         lambda: event_log(args[0] if args else "System", args[1] if len(args)>1 else "Error", int(args[2]) if len(args)>2 else 10),
        "datetime_config":   lambda: datetime_config(args[0], args[1] if len(args)>1 else "", " ".join(args[2:]) if len(args)>2 else ""),
        "ui_auto":           lambda: ui_auto(args[0], args[1] if len(args)>1 else "", args[2] if len(args)>2 else "", " ".join(args[3:]) if len(args)>3 else ""),
        "macro":             lambda: macro(args[0], args[1] if len(args)>1 else "", int(args[2]) if len(args)>2 else 1, float(args[3]) if len(args)>3 else 10.0),
        "color_pick":        lambda: color_pick(int(args[0]), int(args[1]), args[2] if len(args)>2 else "pick", int(args[3]) if len(args)>3 else 100, int(args[4]) if len(args)>4 else 100),
        "webcam":            lambda: webcam(args[0], int(args[1]) if len(args)>1 and args[1].isdigit() else 0, float(args[2]) if len(args)>2 else 5.0, args[3] if len(args)>3 else ""),
        "multi_monitor":     lambda: multi_monitor(args[0], int(args[1]) if len(args)>1 else 1, args[2] if len(args)>2 else ""),
        "printer":           lambda: printer(args[0], args[1] if len(args)>1 else "", args[2] if len(args)>2 else ""),
        "wifi":              lambda: wifi(args[0], args[1] if len(args)>1 else "", args[2] if len(args)>2 else ""),
        "proxy":             lambda: proxy(args[0], args[1] if len(args)>1 else ""),
        "lock_screen":       lambda: lock_screen(args[0] if args else "lock"),
        "defender":          lambda: defender(args[0], args[1] if len(args)>1 else ""),
        "vision_loop":       lambda: vision_loop(" ".join(args), 20, 3.0, 120.0),
        "alert_monitor":     lambda: alert_monitor(args[0], args[1] if len(args)>1 else "80", args[2] if len(args)>2 else "", int(args[3]) if len(args)>3 else 30, int(args[4]) if len(args)>4 else 3600),
        "interval_schedule": lambda: interval_schedule(" ".join(args[:-2]) if len(args)>2 else args[0], float(args[-2]) if len(args)>1 else 60.0, int(args[-1]) if len(args)>2 else 0),
        "wait_for_text":     lambda: wait_for_text(" ".join(args[:-1]) if len(args)>1 else args[0], float(args[-1]) if len(args)>1 and args[-1].replace(".","").isdigit() else 60.0),
        "data_process":      lambda: data_process(args[0], args[1] if len(args)>1 else "", args[2] if len(args)>2 else "", " ".join(args[3:]) if len(args)>3 else ""),
        "wake_on_lan":       lambda: wake_on_lan(args[0], args[1] if len(args)>1 else "255.255.255.255", int(args[2]) if len(args)>2 else 9),
        "clipboard_history": lambda: execute_clipboard_history(args[0] if args else "list", int(args[1]) if len(args)>1 else 0),
        "file_watcher":      lambda: file_watcher(args[0], args[1] if len(args)>1 else "all", args[2] if len(args)>2 else "", float(args[3]) if len(args)>3 else 3600.0),
        "pixel_watch":       lambda: pixel_watch(int(args[0]), int(args[1]), int(args[2]) if len(args)>2 else 10, float(args[3]) if len(args)>3 else 1.0, float(args[4]) if len(args)>4 else 60.0, args[5] if len(args)>5 else ""),
        "object_detect":     lambda: object_detect(" ".join(args[:-1]) if len(args)>1 else args[0], args[-1] if len(args)>1 and args[-1] in ("click","double_click","find") else "find"),
        "mouse_record":      lambda: mouse_record(args[0], args[1] if len(args)>1 else "", float(args[2]) if len(args)>2 else 10.0, int(args[3]) if len(args)>3 else 1, float(args[4]) if len(args)>4 else 1.0),
        "adb":               lambda: adb(args[0], int(args[1]) if len(args)>1 and args[1].isdigit() else 0, int(args[2]) if len(args)>2 and args[2].isdigit() else 0, 0, 0, " ".join(args[3:]) if len(args)>3 else ""),
        "wifi_hotspot":      lambda: wifi_hotspot(args[0], args[1] if len(args)>1 else "", args[2] if len(args)>2 else ""),
        "onedrive":          lambda: onedrive(args[0], args[1] if len(args)>1 else "", args[2] if len(args)>2 else ""),
        "ftp":               lambda: ftp(args[0], args[1] if len(args)>1 else "", args[2] if len(args)>2 else "", args[3] if len(args)>3 else "", args[4] if len(args)>4 else "", args[5] if len(args)>5 else ""),
        "wsl":               lambda: wsl(args[0], args[1] if len(args)>1 else "", " ".join(args[2:]) if len(args)>2 else ""),
        "hyperv":            lambda: hyperv(args[0], args[1] if len(args)>1 else "", args[2] if len(args)>2 else ""),
        "file_diff":         lambda: file_diff(args[0], args[1], args[2] if len(args)>2 else "", args[3] if len(args)>3 else "unified"),
        "screen_live":       lambda: screen_live(float(args[0]) if args else 0.5, float(args[1]) if len(args)>1 else 60.0, int(args[2]) if len(args)>2 else 50),
        # ── 影片生成 ──────────────────────────────────────────────────
        # 用法：video_gen <mode> [output] [key=value ...]
        # mode: slideshow / text_video / tts_video / screen_record
        # 例：video_gen text_video "" text="你好世界" duration=5
        "ai_video":          lambda: ai_video(
            " ".join(args[:1]) if args else "",
            args[1] if len(args) > 1 else "replicate",
            args[2] if len(args) > 2 else "",
            args[3] if len(args) > 3 else "",
            float(args[4]) if len(args) > 4 else 5,
            args[5] if len(args) > 5 else ""
        ),
        "video_gen":         lambda: video_gen(
            args[0] if args else "text_video",
            args[1] if len(args) > 1 else "",
            **dict(
                kv.split("=", 1) for kv in args[2:]
                if "=" in kv
            )
        ),
        # 缺口1
        "email_trigger":     lambda: email_trigger(args[0] if args else "check", *args[1:]),
        "file_trigger":      lambda: file_trigger(args[0] if args else "", args[1] if len(args)>1 else "any", args[2] if len(args)>2 else "notify"),
        "webhook_server":    lambda: webhook_server(args[0] if args else "status", *args[1:]),
        # 缺口2
        "com_auto":          lambda: com_auto(args[0] if args else "excel", args[1] if len(args)>1 else "list_sheets", *args[2:]),
        "dialog_auto":       lambda: dialog_auto(args[0] if args else "list_dialogs", *args[1:]),
        "ime_switch":        lambda: ime_switch(args[0] if args else "status"),
        # 缺口3
        "wake_word":         lambda: wake_word(args[0] if args else "listen_once", *args[1:]),
        "sound_detect":      lambda: sound_detect(args[0] if args else "volume_level", *args[1:]),
        "face_recognize":    lambda: face_recognize(args[0] if args else "detect", *args[1:]),
        # 缺口4
        "http_server":       lambda: http_server(args[0] if args else "status", *args[1:]),
        "lan_scan":          lambda: lan_scan(args[0] if args else "get_local_ip", *args[1:]),
        "serial_port":       lambda: serial_port(args[0] if args else "list", *args[1:]),
        "mqtt":              lambda: mqtt(args[0] if args else "test_connect", args[1] if len(args)>1 else "localhost", *args[2:]),
        # 缺口5
        "doc_ai":            lambda: doc_ai(args[0] if args else "summarize", *args[1:]),
        "web_monitor":       lambda: web_monitor(args[0] if args else "check_once", args[1] if len(args)>1 else "", *args[2:]),
        "audio_transcribe":  lambda: audio_transcribe(args[0] if args else "transcribe_mic", *args[1:]),
        # 投資技能
        "get_institutional":     lambda: get_institutional(args[0] if args else "", args[1] if len(args)>1 else ""),
        "get_sector":            lambda: get_sector(args[0] if args else "us"),
        "get_commodity":         lambda: get_commodity(args[0] if args else "all"),
        "get_bond_yield":        lambda: get_bond_yield(),
        "get_dividend_calendar": lambda: get_dividend_calendar(args[0]),
        "stock_screener":        lambda: stock_screener(args[0], args[1] if len(args)>1 else "us"),
        "get_margin_trading":    lambda: get_margin_trading(args[0], args[1] if len(args)>1 else ""),
        "get_options":           lambda: get_options(args[0], args[1] if len(args)>1 else ""),
        "get_futures":           lambda: get_futures(args[0] if args else "all"),
        "get_ipo":               lambda: get_ipo(args[0] if args else 10),
        "backtest":              lambda: backtest(args[0], args[1] if len(args)>1 else "ma_cross", args[2] if len(args)>2 else "2y"),
        # 中國大陸技能
        "get_ashare":            lambda: get_ashare(args[0], args[1] if len(args)>1 else "1mo"),
        "get_cn_news":           lambda: get_cn_news(args[0] if args else "all", args[1] if len(args)>1 else 5),
        "china_search":          lambda: china_search(args[0], args[1] if len(args)>1 else "其他", args[2] if len(args)>2 else 6),
        # 進階投資技能
        "get_global_market":     lambda: get_global_market(),
        "get_economic_calendar": lambda: get_economic_calendar(args[0] if args else 10),
        "get_earnings_calendar": lambda: get_earnings_calendar(args[0] if args else 7),
        "get_analyst_ratings":   lambda: get_analyst_ratings(args[0]),
        "get_short_interest":    lambda: get_short_interest(args[0]),
        "get_correlation":       lambda: get_correlation(args, args[-1] if args[-1] in ["3mo","6mo","1y","2y"] else "1y"),
        "get_risk_metrics":      lambda: get_risk_metrics(args[0], args[1] if len(args)>1 else "1y"),
        "get_money_flow":        lambda: get_money_flow(args[0]),
        "get_concept_stocks":    lambda: get_concept_stocks(args[0]),
        "get_crypto_depth":      lambda: get_crypto_depth(args[0] if args else "bitcoin"),
        "drip_calculator":       lambda: drip_calculator(args[0], float(args[1]) if len(args)>1 else 100, int(args[2]) if len(args)>2 else 10, float(args[3]) if len(args)>3 else 0),
        "get_forex_chart":       lambda: get_forex_chart(args[0], args[1] if len(args)>1 else "3mo"),
        "get_warrant":           lambda: get_warrant(args[0]),
        "get_portfolio_risk":    lambda: get_portfolio_risk([{"symbol": s, "weight": 1/len(args)} for s in args], "1y"),
        # 理財規劃技能
        "retirement_calculator": lambda: retirement_calculator(*args),
        "loan_calculator":       lambda: loan_calculator(*args),
        "compound_calculator":   lambda: compound_calculator(*args),
        "asset_allocation":      lambda: asset_allocation(*args),
        "tw_tax_calculator":     lambda: tw_tax_calculator(*args),
        "currency_converter":    lambda: currency_converter(*args),
        "get_fund":              lambda: get_fund(args[0]),
        "get_reits":             lambda: get_reits(args[0]),
        "inflation_adjusted":    lambda: inflation_adjusted(*args),
        "defi_calculator":       lambda: defi_calculator(*args),
        "gold_calculator":       lambda: gold_calculator(*args),
        "forex_deposit":         lambda: forex_deposit(*args),
        "financial_health":      lambda: financial_health(*args),
        # 研究統整與觀點技能
        "deep_research":         lambda: deep_research(args[0], args[1] if len(args)>1 else "zh-tw", args[2] if len(args)>2 else 5),
        "fact_check":            lambda: fact_check(" ".join(args)),
        "timeline_events":       lambda: timeline_events(args[0], args[1] if len(args)>1 else "zh-tw"),
        "sentiment_scan":        lambda: sentiment_scan(args[0], args[1] if len(args)>1 else "zh-tw"),
        "compare_analysis":      lambda: compare_analysis(args[0], args[1] if len(args)>1 else ""),
        "pros_cons_analysis":    lambda: pros_cons_analysis(args[0], args[1] if len(args)>1 else "", args[2] if len(args)>2 else "zh-tw"),
        "research_report":       lambda: research_report(args[0], args[1] if len(args)>1 else "一般研究", args[2] if len(args)>2 else "zh-tw"),
        "opinion_writer":        lambda: opinion_writer(args[0], args[1] if len(args)>1 else "中立", args[2] if len(args)>2 else "正式"),
        "trend_forecast":        lambda: trend_forecast(args[0], args[1] if len(args)>1 else "全部", args[2] if len(args)>2 else "zh-tw"),
        "debate_simulator":      lambda: debate_simulator(" ".join(args)),
        # 全域研究技能
        "academic_search":       lambda: academic_search(args[0], args[1] if len(args)>1 else "", args[2] if len(args)>2 else "en"),
        "health_research":       lambda: health_research(" ".join(args)),
        "law_research":          lambda: law_research(args[0], args[1] if len(args)>1 else "台灣", args[2] if len(args)>2 else "zh-tw"),
        "person_research":       lambda: person_research(args[0], args[1] if len(args)>1 else ""),
        "company_research":      lambda: company_research(args[0]),
        "product_review":        lambda: product_review(args[0], args[1] if len(args)>1 else ""),
        "travel_research":       lambda: travel_research(args[0], args[1] if len(args)>1 else None, args[2] if len(args)>2 else ""),
        "job_market":            lambda: job_market(args[0], args[1] if len(args)>1 else "台灣"),
        "impact_analysis":       lambda: impact_analysis(args[0], args[1] if len(args)>1 else ""),
        "scenario_planning":     lambda: scenario_planning(args[0], args[1] if len(args)>1 else ""),
        "decision_helper":       lambda: decision_helper(args[0], args[1] if len(args)>1 else "", args[2] if len(args)>2 else ""),
        "devil_advocate":        lambda: devil_advocate(" ".join(args)),
        # 統整與觀點強化技能
        "summary_writer":        lambda: summary_writer(args[0], args[1] if len(args)>1 else 7),
        "key_insights":          lambda: key_insights(args[0], args[1] if len(args)>1 else 5),
        "bias_detector":         lambda: bias_detector(" ".join(args)),
        "second_opinion":        lambda: second_opinion(args[0], args[1] if len(args)>1 else ""),
        "brainstorm":            lambda: brainstorm(args[0], args[1] if len(args)>1 else 8, args[2] if len(args)>2 else "實用"),
        "benchmark_analysis":    lambda: benchmark_analysis(args[0], args[1] if len(args)>1 else ""),
        # 觀點表達技能
        "steel_man":             lambda: steel_man(args[0], args[1] if len(args)>1 else ""),
        "socratic_questioning":  lambda: socratic_questioning(args[0], args[1] if len(args)>1 else 5),
        "analogy_maker":         lambda: analogy_maker(args[0], args[1] if len(args)>1 else "一般大眾", args[2] if len(args)>2 else 3),
        "narrative_builder":     lambda: narrative_builder(args[0], args[1] if len(args)>1 else "", args[2] if len(args)>2 else ""),
        "critique_writer":       lambda: critique_writer(args[0], args[1] if len(args)>1 else "觀點"),
        "position_statement":    lambda: position_statement(args[0], args[1] if len(args)>1 else "支持"),
        # 桌面視覺控制技能
        "ocr_click":             lambda: ocr_click(args[0], args[1] if len(args)>1 else 1, args[2] if len(args)>2 else "click"),
        "vision_locate":         lambda: vision_locate(args[0], args[1] if len(args)>1 else 1, args[2] if len(args)>2 else "click"),
        "screen_workflow":       lambda: screen_workflow(args[0]),
        "app_navigator":         lambda: app_navigator(args[0], args[1] if len(args)>1 else "", args[2] if len(args)>2 else "", args[3] if len(args)>3 else 1),
        "wait_and_click":        lambda: wait_and_click(args[0], args[1] if len(args)>1 else 15, args[2] if len(args)>2 else 1),
        "drag_drop":             lambda: drag_drop(args[0] if args else None, args[1] if len(args)>1 else None, args[2] if len(args)>2 else None, args[3] if len(args)>3 else None),
        # ── 新增 116 個工具 ──────────────────────────────────
        "analyze_pdf":           lambda: analyze_pdf_tool(args[0], args[1] if len(args)>1 else 4000),
        "audio_process":         lambda: audio_process_tool(args[0], args[1], args[2] if len(args)>2 else "", args[3] if len(args)>3 else 0, args[4] if len(args)>4 else 0),
        "auto_skill":            lambda: auto_skill_tool(args[0], args[1] if len(args)>1 else "", args[2] if len(args)>2 else "", " ".join(args[3:]) if len(args)>3 else ""),
        "auto_trade":            lambda: auto_trade_tool(args[0], args[1] if len(args)>1 else "", args[2] if len(args)>2 else 0, args[3] if len(args)>3 else 0, args[4] if len(args)>4 else "market"),
        "automation":            lambda: automation_tool(args[0], args[1] if len(args)>1 else "", args[2] if len(args)>2 else "", " ".join(args[3:]) if len(args)>3 else ""),
        "barcode":               lambda: barcode_tool(args[0] if args else ""),
        "bluetooth":             lambda: bluetooth_tool(args[0], args[1] if len(args)>1 else ""),
        "browser_advanced":      lambda: browser_advanced_tool(args[0], args[1] if len(args)>1 else "", args[2] if len(args)>2 else "", args[3] if len(args)>3 else "", args[4] if len(args)>4 else 0, args[5] if len(args)>5 else 30.0, args[6] if len(args)>6 else ""),
        "browser_control":       lambda: browser_control_tool(args[0], args[1] if len(args)>1 else "", args[2] if len(args)>2 else "", " ".join(args[3:]) if len(args)>3 else ""),
        "calendar":              lambda: calendar_tool(args[0], args[1] if len(args)>1 else 7, args[2] if len(args)>2 else "", args[3] if len(args)>3 else "", args[4] if len(args)>4 else "", " ".join(args[5:]) if len(args)>5 else ""),
        "clipboard":             lambda: clipboard_tool(args[0], " ".join(args[1:]) if len(args)>1 else ""),
        "clipboard_image":       lambda: clipboard_image_tool(args[0], args[1] if len(args)>1 else ""),
        "cloud_storage":         lambda: cloud_storage_tool(args[0], args[1], args[2] if len(args)>2 else "root"),
        "compare_stocks":        lambda: compare_stocks_tool(args[0], args[1] if len(args)>1 else "all"),
        "database":              lambda: database_tool(args[0], args[1], " ".join(args[2:]) if len(args)>2 else ""),
        "ddg_search":            lambda: ddg_search_tool(" ".join(args[:1] if len(args)==1 else [args[0]]), args[1] if len(args)>1 else "zh-tw", args[2] if len(args)>2 else 5),
        "desktop_control":       lambda: desktop_control_tool(args[0], args[1] if len(args)>1 else None, args[2] if len(args)>2 else None, args[3] if len(args)>3 else None, args[4] if len(args)>4 else None),
        "device_manager":        lambda: device_manager_tool(args[0], args[1] if len(args)>1 else "", args[2] if len(args)>2 else ""),
        "disk_backup":           lambda: disk_backup_tool(args[0], args[1] if len(args)>1 else "", args[2] if len(args)>2 else ""),
        "display":               lambda: display_tool(args[0], args[1] if len(args)>1 else None),
        "docker":                lambda: docker_tool(args[0], args[1] if len(args)>1 else ""),
        "document_control":      lambda: document_control_tool(args[0], args[1], " ".join(args[2:]) if len(args)>2 else ""),
        "download_file":         lambda: download_file_tool(args[0], args[1] if len(args)>1 else ""),
        "dropbox":               lambda: dropbox_tool(args[0], args[1], args[2], args[3] if len(args)>3 else ""),
        "email_control":         lambda: email_control_tool(args[0], args[1], args[2], args[3] if len(args)>3 else "INBOX", args[4] if len(args)>4 else 5),
        "emotion_detect":        lambda: emotion_detect_tool(args[0], " ".join(args[1:]) if len(args)>1 else "", ""),
        "encrypt_file":          lambda: encrypt_file_tool(args[0], args[1], args[2]),
        "env_var":               lambda: env_var_tool(args[0], args[1] if len(args)>1 else "", args[2] if len(args)>2 else "", args[3] if len(args)>3 else "false"),
        "file_system":           lambda: file_system_tool(args[0], args[1] if len(args)>1 else "", args[2] if len(args)>2 else "", " ".join(args[3:]) if len(args)>3 else ""),
        "file_tools":            lambda: file_tools_tool(args[0], args[1], args[2] if len(args)>2 else "", args[3] if len(args)>3 else "", args[4] if len(args)>4 else "", args[5] if len(args)>5 else ""),
        "file_transfer":         lambda: file_transfer_tool(args[0], args[1], args[2] if len(args)>2 else ""),
        "find_image_on_screen":  lambda: find_image_on_screen_tool(args[0], args[1] if len(args)>1 else 0.8),
        "generate_image":        lambda: generate_image_tool(args[0], args[1] if len(args)>1 else 512, args[2] if len(args)>2 else 512, " ".join(args[3:]) if len(args)>3 else ""),
        "get_candlestick_chart": lambda: get_candlestick_chart_tool(args[0], args[1] if len(args)>1 else "3mo"),
        "get_crypto":            lambda: get_crypto_tool(args[0], args[1] if len(args)>1 else "usd"),
        "get_earnings":          lambda: get_earnings_tool(args[0]),
        "get_etf":               lambda: get_etf_tool(args[0]),
        "get_finance_news":      lambda: get_finance_news_tool(args[0] if args else "all", args[1] if len(args)>1 else 5),
        "get_forex":             lambda: get_forex_tool(args[0]),
        "get_fundamentals":      lambda: get_fundamentals_tool(args[0]),
        "get_macro":             lambda: get_macro_tool(args[0]),
        "get_market_sentiment":  lambda: get_market_sentiment_tool(),
        "get_stock":             lambda: get_stock_tool(args[0], args[1] if len(args)>1 else "1mo"),
        "get_stock_advanced":    lambda: get_stock_advanced_tool(args[0], args[1] if len(args)>1 else "macd,bb,kd"),
        "get_weather":           lambda: get_weather_tool(args[0]),
        "git":                   lambda: git_tool(args[0], args[1] if len(args)>1 else ".", args[2] if len(args)>2 else "", args[3] if len(args)>3 else "master"),
        "goal_manager":          lambda: goal_manager_tool(args[0], args[1] if len(args)>1 else "", args[2] if len(args)>2 else "", args[3] if len(args)>3 else "", args[4] if len(args)>4 else "normal"),
        "google_trends":         lambda: google_trends_tool(args[0], args[1] if len(args)>1 else "today 3-m", args[2] if len(args)>2 else "TW"),
        "hardware":              lambda: hardware_tool(),
        "image_edit":            lambda: image_edit_tool(args[0], args[1], *args[2:]),
        "image_tools":           lambda: image_tools_tool(args[0], args[1] if len(args)>1 else "", args[2] if len(args)>2 else 75, args[3] if len(args)>3 else 0, args[4] if len(args)>4 else 0),
        "knowledge_base":        lambda: knowledge_base_tool(args[0], " ".join(args[1:]) if len(args)>1 else ""),
        "long_term_memory":      lambda: long_term_memory_tool(args[0], args[1] if len(args)>1 else "0", " ".join(args[2:]) if len(args)>2 else "", args[3] if len(args)>3 else ""),
        "lookup":                lambda: lookup_tool(args[0], args[1] if len(args)>1 else "", args[2] if len(args)>2 else 1.0, args[3] if len(args)>3 else "USD", args[4] if len(args)>4 else "TWD"),
        "manage_schedule":       lambda: manage_schedule_tool(args[0], args[1] if len(args)>1 else "", args[2] if len(args)>2 else "", args[3] if len(args)>3 else ""),
        "media":                 lambda: media_tool(args[0], args[1] if len(args)>1 else ""),
        "monitor_config":        lambda: monitor_config_tool(),
        "multi_deploy":          lambda: multi_deploy_tool(args[0], args[1] if len(args)>1 else "", args[2] if len(args)>2 else "", args[3] if len(args)>3 else "", args[4] if len(args)>4 else "/tmp/niu_bot"),
        "multi_perspective":     lambda: multi_perspective_tool(" ".join(args), "zh-tw"),
        "network_config":        lambda: network_config_tool(args[0], args[1] if len(args)>1 else "", args[2] if len(args)>2 else "", args[3] if len(args)>3 else "", args[4] if len(args)>4 else ""),
        "network_diag":          lambda: network_diag_tool(args[0], args[1], args[2] if len(args)>2 else "22,80,443,3306,3389,8080"),
        "news_monitor":          lambda: news_monitor_tool(args[0], args[1] if len(args)>1 else "", args[2] if len(args)>2 else 300, args[3] if len(args)>3 else 3600),
        "news_search":           lambda: news_search_tool(" ".join(args[:1] if len(args)==1 else [args[0]]), args[1] if len(args)>1 else "zh-TW", args[2] if len(args)>2 else 6),
        "nlp":                   lambda: nlp_tool(args[0], " ".join(args[1:])),
        "osint_search":          lambda: osint_search_tool(args[0], args[1] if len(args)>1 else "", args[2] if len(args)>2 else "", args[3] if len(args)>3 else 10),
        "password_mgr":          lambda: password_mgr_tool(args[0], args[1], args[2], args[3] if len(args)>3 else "", args[4] if len(args)>4 else ""),
        "pdf_edit":              lambda: pdf_edit_tool(args[0], args[1] if len(args)>1 else "", args[2] if len(args)>2 else "", args[3] if len(args)>3 else "", " ".join(args[4:]) if len(args)>4 else ""),
        "pdf_image":             lambda: pdf_image_tool(args[0], args[1] if len(args)>1 else "", args[2] if len(args)>2 else 150),
        "pentest":               lambda: pentest_tool(args[0], args[1] if len(args)>1 else "", args[2] if len(args)>2 else "1-1000", args[3] if len(args)>3 else 2),
        "portfolio":             lambda: portfolio_tool(args[0], args[1] if len(args)>1 else 0, args[2] if len(args)>2 else "", args[3] if len(args)>3 else 0, args[4] if len(args)>4 else 0),
        "power_control":         lambda: power_control_tool(args[0]),
        "pptx_control":          lambda: pptx_control_tool(args[0], args[1], " ".join(args[2:]) if len(args)>2 else ""),
        "proactive_alert":       lambda: proactive_alert_tool(args[0], args[1] if len(args)>1 else "", args[2] if len(args)>2 else "", args[3] if len(args)>3 else "", args[4] if len(args)>4 else "", args[5] if len(args)>5 else 60),
        "ptt_search":            lambda: ptt_search_tool(args[0], args[1] if len(args)>1 else "Gossiping", args[2] if len(args)>2 else 5),
        "push_notify":           lambda: push_notify_tool(args[0], " ".join(args[1:-1]), args[-1]),
        "qr_code":               lambda: qr_code_tool(args[0], args[1] if len(args)>1 else "", args[2] if len(args)>2 else "", args[3] if len(args)>3 else 30.0),
        "read_screen":           lambda: read_screen_tool(" ".join(args) if args else "描述螢幕上有什麼", args[-1] if args and args[-1].isdigit() else 1),
        "read_webpage":          lambda: read_webpage_tool(args[0], args[1] if len(args)>1 else 3000),
        "registry":              lambda: registry_tool(args[0], args[1], args[2] if len(args)>2 else "", " ".join(args[3:]) if len(args)>3 else ""),
        "reminder":              lambda: reminder_tool(args[0], " ".join(args[1:])),
        "report":                lambda: report_tool(args[0], args[1], args[2] if len(args)>2 else ""),
        "restore_point":         lambda: restore_point_tool(args[0], " ".join(args[1:]) if len(args)>1 else ""),
        "run_code":              lambda: run_code_tool(args[0], " ".join(args[1:])),
        "screen_vision":         lambda: screen_vision_tool(" ".join(args) if args else "請描述這個畫面上有什麼"),
        "scroll_at":             lambda: scroll_at_tool(args[0] if args else "down", args[1] if len(args)>1 else 3, args[2] if len(args)>2 else None, args[3] if len(args)>3 else None, args[4] if len(args)>4 else 1, " ".join(args[5:]) if len(args)>5 else ""),
        "self_benchmark":        lambda: self_benchmark_tool(args[0]),
        "send_email":            lambda: send_email_tool(args[0], args[1], " ".join(args[2:])),
        "send_voice":            lambda: send_voice_tool(" ".join(args[:-1]) if len(args)>1 else " ".join(args), args[-1] if len(args)>1 and "Neural" in args[-1] else "zh-CN-YunxiNeural"),
        "smart_home":            lambda: smart_home_tool(args[0], args[1] if len(args)>1 else "", args[2] if len(args)>2 else "", args[3] if len(args)>3 else "", args[4] if len(args)>4 else ""),
        "software":              lambda: software_tool(args[0], " ".join(args[1:]) if len(args)>1 else ""),
        "ssh_sftp":              lambda: ssh_sftp_tool(args[0], args[1], args[2], args[3], " ".join(args[4:]) if len(args)>4 else ""),
        "startup":               lambda: startup_tool(args[0], args[1] if len(args)>1 else "", " ".join(args[2:]) if len(args)>2 else ""),
        "system_monitor":        lambda: system_monitor_tool(args[0], args[1] if len(args)>1 else ""),
        "system_tools":          lambda: system_tools_tool(args[0]),
        "think_as":              lambda: think_as_tool(args[0], " ".join(args[1:]) if len(args)>1 else ""),
        "threat_intel":          lambda: threat_intel_tool(args[0], args[1] if len(args)>1 else "", args[2] if len(args)>2 else ""),
        "todo_list":             lambda: todo_list_tool(args[0], " ".join(args[1:]) if len(args)>1 else "", 0),
        "tts_advanced":          lambda: tts_advanced_tool(args[0], " ".join(args[1:]) if len(args)>1 else "", args[-1] if len(args)>1 and "Neural" in args[-1] else "zh-CN-YunxiNeural"),
        "user_account":          lambda: user_account_tool(args[0], args[1] if len(args)>1 else "", args[2] if len(args)>2 else ""),
        "video_gif":             lambda: video_gif_tool(args[0], args[1] if len(args)>1 else 0, args[2] if len(args)>2 else 5.0, args[3] if len(args)>3 else "", args[4] if len(args)>4 else 10),
        "video_process":         lambda: video_process_tool(args[0], args[1], args[2] if len(args)>2 else 0, args[3] if len(args)>3 else 0, args[4] if len(args)>4 else 0, args[5] if len(args)>5 else ""),
        "virtual_desktop":       lambda: virtual_desktop_tool(args[0]),
        "voice_cmd":             lambda: voice_cmd_tool(args[0], args[1] if len(args)>1 else 300.0, args[2] if len(args)>2 else "zh-TW"),
        "voice_id":              lambda: voice_id_tool(args[0], args[1] if len(args)>1 else "", args[2] if len(args)>2 else "", args[3] if len(args)>3 else 5),
        "volume":                lambda: volume_tool(args[0], args[1] if len(args)>1 else None),
        "vpn":                   lambda: vpn_tool(args[0], args[1] if len(args)>1 else "", args[2] if len(args)>2 else "", args[3] if len(args)>3 else ""),
        "wait_seconds":          lambda: wait_seconds_tool(args[0]),
        "web_scrape":            lambda: web_scrape_tool(args[0], args[1] if len(args)>1 else "", args[2] if len(args)>2 else "body", args[3] if len(args)>3 else 2.0, args[4] if len(args)>4 else "full"),
        "webpage_shot":          lambda: webpage_shot_tool(args[0], args[1], args[2] if len(args)>2 else "body", args[3] if len(args)>3 else 60.0, args[4] if len(args)>4 else 3600.0),
        "wikipedia_search":      lambda: wikipedia_search_tool(args[0], args[1] if len(args)>1 else "zh"),
        "win_notify_relay":      lambda: win_notify_relay_tool(args[0], args[1] if len(args)>1 else 3600.0, args[2] if len(args)>2 else ""),
        "window_control":        lambda: window_control_tool(args[0], " ".join(args[1:]) if len(args)>1 else ""),
        "window_manager":        lambda: window_manager_tool(args[0] if args else "list", " ".join(args[1:]) if len(args)>1 else ""),
        "windows_update":        lambda: windows_update_tool(args[0]),
        "workflow":              lambda: workflow_tool(args[0], args[1] if len(args)>1 else "", " ".join(args[2:]) if len(args)>2 else ""),
        "youtube_summary":       lambda: youtube_summary_tool(args[0]),
        "tg_auto_reply":         lambda: tg_auto_reply_tool(args[0] if args else "start", args[1] if len(args)>1 else "", args[2] if len(args)>2 else ""),
        "line_send_msg":         lambda: line_send_msg_tool(args[0] if args else "", " ".join(args[1:]) if len(args)>1 else ""),
        # ── 29 個 execute_* 函數（與 bot.py 同步）──────────────
        "execute_ai_plan":       lambda: print(execute_ai_plan(" ".join(args))),
        "execute_api_call":      lambda: print(execute_api_call(args[0], args[1], args[2] if len(args)>2 else "{}", args[3] if len(args)>3 else "{}")),
        "execute_chrome_bookmarks": lambda: print(execute_chrome_bookmarks()),
        "execute_defender":      lambda: print(execute_defender(args[0], args[1] if len(args)>1 else "")),
        "execute_disk_clean":    lambda: print(execute_disk_clean(args[0] if args else "list")),
        "execute_drag":          lambda: print(execute_drag(args[0], args[1], args[2], args[3], float(args[4]) if len(args)>4 else 0.5)),
        "execute_firewall":      lambda: print(execute_firewall(args[0], args[1] if len(args)>1 else "", int(args[2]) if len(args)>2 else None, args[3] if len(args)>3 else "TCP", args[4] if len(args)>4 else "Inbound")),
        "execute_font_list":     lambda: print(execute_font_list(args[0] if args else "")),
        "execute_ftp":           lambda: print(execute_ftp(args[0], args[1] if len(args)>1 else "", args[2] if len(args)>2 else "", args[3] if len(args)>3 else "", args[4] if len(args)>4 else "", args[5] if len(args)>5 else "", int(args[6]) if len(args)>6 else 21)),
        "execute_hotkey":        lambda: print(execute_hotkey(args[0])),
        "execute_hyperv":        lambda: print(execute_hyperv(args[0], args[1] if len(args)>1 else "", args[2] if len(args)>2 else "")),
        "execute_interval_schedule": lambda: print(execute_interval_schedule(args[0], args[1] if len(args)>1 else "", args[2] if len(args)>2 else "", float(args[3]) if len(args)>3 else 60.0, int(args[4]) if len(args)>4 else 0, float(args[5]) if len(args)>5 else 0.0)),
        "execute_lock_screen":   lambda: print(execute_lock_screen(args[0] if args else "lock")),
        "execute_net_share":     lambda: print(execute_net_share(args[0], args[1] if len(args)>1 else "", args[2] if len(args)>2 else "Z:", args[3] if len(args)>3 else "", args[4] if len(args)>4 else "")),
        "execute_object_detect": lambda: print(execute_object_detect(" ".join(args[:-1]) if len(args)>1 else args[0], args[-1] if len(args)>1 and args[-1] in ("click","double_click","find") else "find")),
        "execute_onedrive":      lambda: print(execute_onedrive(args[0], args[1] if len(args)>1 else "", args[2] if len(args)>2 else "")),
        "execute_proxy":         lambda: print(execute_proxy(args[0], args[1] if len(args)>1 else "")),
        "execute_rdp_connect":   lambda: print(execute_rdp_connect(args[0], args[1] if len(args)>1 else "", int(args[2]) if len(args)>2 else 1280, int(args[3]) if len(args)>3 else 720)),
        "execute_right_menu":    lambda: print(execute_right_menu(args[0], args[1], args[2] if len(args)>2 else "")),
        "execute_screen_live":   lambda: print(execute_screen_live(args[0] if args else "start", float(args[1]) if len(args)>1 else 0.5, float(args[2]) if len(args)>2 else 60.0, int(args[3]) if len(args)>3 else 50)),
        "execute_screen_stream": lambda: print(f"已擷取 {len(execute_screen_stream(int(args[0]) if args else 10, int(args[1]) if len(args)>1 else 2))} 張截圖"),
        "execute_screen_watch":  lambda: print(execute_screen_watch(args[0], args[1], float(args[2]) if len(args)>2 else 60)),
        "execute_sysres_chart":  lambda: print(execute_sysres_chart(int(args[0]) if args else 10)),
        "execute_usb_list":      lambda: print(execute_usb_list()),
        "execute_wake_on_lan":   lambda: print(execute_wake_on_lan(args[0], args[1] if len(args)>1 else "255.255.255.255", int(args[2]) if len(args)>2 else 9)),
        "execute_wifi":          lambda: print(execute_wifi(args[0], args[1] if len(args)>1 else "", args[2] if len(args)>2 else "")),
        "execute_wifi_hotspot":  lambda: print(execute_wifi_hotspot(args[0], args[1] if len(args)>1 else "", args[2] if len(args)>2 else "")),
        "execute_win_service":   lambda: print(execute_win_service(args[0], args[1] if len(args)>1 else "")),
        "execute_wsl":           lambda: print(execute_wsl(args[0], args[1] if len(args)>1 else "", " ".join(args[2:]) if len(args)>2 else "")),
    }

    if tool not in tools:
        print(f"未知工具：{tool}")
        print(__doc__)
        sys.exit(1)

    try:
        tools[tool]()
    except Exception as e:
        print(f"執行失敗：{e}")
        sys.exit(1)
