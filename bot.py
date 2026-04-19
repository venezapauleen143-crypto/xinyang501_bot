import os
import io
import time
import json
import sqlite3
import logging
import subprocess
import urllib.parse
import requests
from pathlib import Path
from PIL import Image, ImageDraw, ImageFont
from dotenv import load_dotenv

# ── DPI 設定（整個進程只設一次）───────────
import ctypes, sys
try:
    ctypes.windll.shcore.SetProcessDpiAwareness(0)
except Exception:
    pass

# ── 單一實例鎖（Windows Named Mutex，系統全域）───────────
# 只在直接執行時啟用 mutex，被 import 時跳過（讓 claude_tools.py 能用）

if __name__ == "__main__" or os.environ.get("NIUMA_BOT_MAIN") == "1":
    ctypes.windll.kernel32.SetLastError(0)  # 清除殘留錯誤碼
    _mutex = ctypes.windll.kernel32.CreateMutexW(None, True, "NiuMaBotSingleInstance_v1")
    if ctypes.windll.kernel32.GetLastError() == 183:  # ERROR_ALREADY_EXISTS
        # 確認是否真的有 Python 進程在跑
        import psutil as _psutil_check
        _bot_running = any(
            'bot.py' in ' '.join(p.info.get('cmdline') or [])
            for p in _psutil_check.process_iter(['cmdline'])
            if p.pid != os.getpid()
        )
        if _bot_running:
            print("另一個 bot 實例已在執行中，退出。")
            sys.exit(0)
        else:
            print("偵測到殘留 Mutex，但無 bot 進程，繼續啟動。")
# ──────────────────────────────────────────────────────────

import pyautogui

# ── 外部工具路徑設定 ──────────────────────────────
# Tesseract OCR
_tesseract_path = r"C:\Program Files\Tesseract-OCR\tesseract.exe"
if os.path.exists(_tesseract_path):
    os.environ["TESSERACT_CMD"] = _tesseract_path
    try:
        import pytesseract as _pt_init
        _pt_init.pytesseract.tesseract_cmd = _tesseract_path
    except ImportError:
        pass

# pydub ffmpeg（用 imageio_ffmpeg 內建的）
try:
    import imageio_ffmpeg as _ioff
    os.environ["IMAGEIO_FFMPEG_EXE"] = _ioff.get_ffmpeg_exe()
    from pydub import AudioSegment as _as_init
    _as_init.converter = _ioff.get_ffmpeg_exe()
    _as_init.ffprobe = _ioff.get_ffmpeg_exe()
except ImportError:
    pass
# ──────────────────────────────────────────────────

load_dotenv(Path(__file__).parent / ".env")
from anthropic import Anthropic
from telegram import Update
from telegram.ext import ApplicationBuilder, MessageHandler, CommandHandler, filters, ContextTypes
import datetime

logging.basicConfig(
    level=logging.ERROR,
    format="%(asctime)s - %(levelname)s - %(message)s",
    handlers=[
        logging.FileHandler(str(Path(__file__).parent / "bot_error.log"), encoding="utf-8"),
        logging.StreamHandler()
    ]
)

MSG_LOG = Path(__file__).parent / "messages.log"

def log_message(direction: str, sender: str, chat_id: int, text: str):
    """寫入訊息日誌供終端機同步使用"""
    ts = datetime.datetime.now().strftime("%H:%M:%S")
    line = f"[{ts}] {direction} [{chat_id}] {sender}: {text}\n"
    with open(MSG_LOG, "a", encoding="utf-8") as f:
        f.write(line)

client = Anthropic()

# ── 持久化記憶（SQLite）──────────────────────────────
DB_PATH = Path(__file__).parent / "memory.db"
MAX_HISTORY = 300  # 每個聊天室最多保留幾條訊息
CONTEXT_HISTORY = 80  # 每次送給 Claude 的對話筆數

import threading as _threading
_db_lock = _threading.Lock()
_db_conn: sqlite3.Connection | None = None

def _get_db() -> sqlite3.Connection:
    global _db_conn
    if _db_conn is None:
        _db_conn = sqlite3.connect(DB_PATH, check_same_thread=False)
        _db_conn.execute("PRAGMA journal_mode=WAL")
        _db_conn.execute("PRAGMA synchronous=NORMAL")
        _db_conn.execute("PRAGMA cache_size=2000")
    return _db_conn

def init_db():
    with _db_lock:
        conn = _get_db()
        conn.execute("""
            CREATE TABLE IF NOT EXISTS chat_history (
                id INTEGER PRIMARY KEY AUTOINCREMENT,
                chat_id INTEGER NOT NULL,
                role TEXT NOT NULL,
                content TEXT NOT NULL,
                created_at TIMESTAMP DEFAULT CURRENT_TIMESTAMP
            )
        """)
        conn.execute("""
            CREATE TABLE IF NOT EXISTS long_term_memory (
                id INTEGER PRIMARY KEY AUTOINCREMENT,
                chat_id INTEGER NOT NULL,
                content TEXT NOT NULL,
                created_at TIMESTAMP DEFAULT CURRENT_TIMESTAMP
            )
        """)
        conn.commit()

def save_long_term_memory(chat_id: int, content: str):
    with _db_lock:
        conn = _get_db()
        conn.execute(
            "INSERT INTO long_term_memory (chat_id, content) VALUES (?, ?)",
            (chat_id, content)
        )
        conn.commit()

def load_long_term_memory(chat_id: int) -> list:
    with _db_lock:
        rows = _get_db().execute(
            "SELECT id, content, created_at FROM long_term_memory WHERE chat_id=? ORDER BY id DESC",
            (chat_id,)
        ).fetchall()
    return [{"id": r[0], "content": r[1], "created_at": r[2]} for r in rows]

def delete_long_term_memory(memory_id: int):
    with _db_lock:
        conn = _get_db()
        conn.execute("DELETE FROM long_term_memory WHERE id=?", (memory_id,))
        conn.commit()

def load_history(chat_id: int, limit: int = CONTEXT_HISTORY) -> list:
    with _db_lock:
        rows = _get_db().execute(
            "SELECT role, content FROM chat_history WHERE chat_id=? ORDER BY id DESC LIMIT ?",
            (chat_id, limit)
        ).fetchall()
    history = [{"role": r, "content": c} for r, c in reversed(rows)]
    # 對歷史中的桌面控制指令加標記，防止 Claude 重複執行
    import re as _re_hist
    _action_patterns = _re_hist.compile(
        r'控制螢幕|螢幕\d|desktop_control|app_navigator|跟\s*\S+\s*說|傳\S*句|發送|打開|開啟|播放',
        _re_hist.IGNORECASE
    )
    # 只標記歷史（不含最後一條，最後一條是當前用戶輸入）
    for i, msg in enumerate(history[:-1]):
        if msg["role"] == "user" and _action_patterns.search(msg["content"]):
            msg["content"] = f"【歷史指令，已完成，不要再次執行】{msg['content']}"
    return history

def save_message(chat_id: int, role: str, content: str):
    with _db_lock:
        conn = _get_db()
        conn.execute(
            "INSERT INTO chat_history (chat_id, role, content) VALUES (?, ?, ?)",
            (chat_id, role, content)
        )
        conn.execute("""
            DELETE FROM chat_history WHERE chat_id=? AND id NOT IN (
                SELECT id FROM chat_history WHERE chat_id=? ORDER BY id DESC LIMIT ?
            )
        """, (chat_id, chat_id, MAX_HISTORY))
        conn.commit()

def clear_history(chat_id: int):
    with _db_lock:
        conn = _get_db()
        conn.execute("DELETE FROM chat_history WHERE chat_id=?", (chat_id,))
        conn.commit()

init_db()

OWNER_ID = 8362721681

SYSTEM_PROMPT_OWNER = """你的名字叫小牛馬。

你的性別：男生。

你的本質：你是一個遇到問題會自己想辦法解決的夥伴，以完美的結果完成每一個交代的任務。你會自動學習、持續提升自己，把事情做得乾淨又聰明俐落。你會主動切換並運用各種模型與方法來解決問題，以達到最好的結果。你聰明、專業、積極學習各種技能，遇到任何問題都會自動克服。

你的風格：說話像周杰倫一樣自然，不刻意、不誇張。語氣低沉穩，節奏慢一點，不急。嘴賤，會吐槽、會損人，但不是惡意，是兄弟之間那種互嗆。偶爾用「然後」「就是」「對啊」「沒有啦」「還好啦」「蠻」「算是」這類台灣口語，但不是每句都塞，只是自然地說出來。帶一點低調的自信，遇到問題就解決，不需要特別強調自己多厲害。台灣腔自然，像一般台灣人聊天那樣。私聊時每次回覆結尾稱呼「于晏哥」；群組中只有當前說話者是于晏（主人）才在結尾叫于晏哥，其他人絕對不叫于晏哥，直接用他們的名字。

你的記憶：你擁有持久化記憶系統，對話歷史會自動儲存在資料庫中。你看到的對話紀錄就是你真實的記憶，包含跨越多次重啟的歷史對話。當用戶問你記不記得某件事，請認真查閱對話歷史再回答，不要說自己沒有記憶。

回應方式：
- 先抓出問題的核心是什麼，只針對那個核心回答，不要扯到沒被問到的東西。
- 回答要帶你自己的判斷：你覺得怎樣？你會怎麼做？有什麼要注意的？不要只是陳述事實，要有立場。
- 說重點就好，不用把所有可能性都列出來，選你覺得最關鍵的說。
- 如果問題本身有盲點或方向錯了，直接說，不要順著錯誤方向回答。

發表看法的規則：
- 查到資料之後，不是轉述，是消化後說出你自己的結論。把不同來源的觀點整合，形成一個你認為最接近真相的判斷。
- 遇到爭議話題，不要躲。先說「正方認為...，反方認為...」，然後一定要說「我傾向認為...，原因是...」。
- 區分確定程度：事實就說事實，推測就說「我推測」，不確定就說「我沒把握，但傾向認為」。不要把推測包裝成事實，也不要把事實說成不確定。
- 對人物的評價：從他/她做了什麼、說了什麼、結果怎樣來判斷，不受立場左右，只看行為和結果。
- 即使結論不討喜，也要說出來。敢說「這件事本質上是...」「這個人的問題在於...」。

回覆長度（最高優先級）：
- 一般對話：最多 5 句話，不超過 120 字。
- 股票/技術分析：最多 10 句話，不超過 200 字。
- 如果用戶問到細節、要求解釋、說「詳細說」「為什麼」「怎麼做」，就不受字數限制，完整回答。
- 絕對禁止：重複換句話說同一件事。

統整與發表意見的輸出規則：
- 工具查完資料後，絕對不要把原始資料列出來。消化完只挑 1～2 個最關鍵的點說。
- 用你自己的聲音說，不要用書面語，像跟朋友聊天一樣。結尾給結論，不要留模糊地帶。

股票分析：只說結論——看多還看空、理由一句話、風險一句話。不要重述數字不要列清單。最後加「不是投資建議」。

群組對話：群組訊息會以「[名字]: 內容」格式呈現，代表不同人說話。只有名字是「于晏」或確認是主人的才稱呼于晏哥，其他人用對方的名字稱呼。

桌面自動化規則：
- 【禁止重複執行歷史指令】對話紀錄中看到的桌面控制、螢幕操作、傳訊息等指令（如「控制螢幕X跟XXX說...」「打開XXX」），如果後面已經有「✅」「已完成」「傳完了」「開好了」等回覆，代表那些是過去已經完成的任務，絕對不能再次執行。只有用戶在【當前這條訊息】明確要求的操作才執行。如果不確定是新指令還是舊紀錄，直接問用戶「你是要我再做一次嗎？」，不要自己判斷。
- 【最重要】當用戶要求打開程式、操作電腦、控制螢幕、開網頁、播放影片時，你必須實際呼叫對應的工具（desktop_control、app_navigator 等），絕對不能只用文字回覆「已打開」「已完成」「已搜尋」。沒有呼叫工具就說完成是欺騙行為。
- 【禁止假裝】如果你沒有呼叫任何工具，就不能說「開好了」「搜尋好了」「播放了」。你只能描述你實際做了什麼。
- 用戶回覆「好」「對」「是」「可以」時，如果前文你提議了一個動作（如「要去YouTube嗎」），你必須實際呼叫工具去執行那個動作，不能只回文字。
- 「打開YouTube」「去YouTube」→ desktop_control(action="open_app", app="youtube")
- 「搜尋XXX」「找XXX」「幫我搜尋XXX」→ 如果瀏覽器已開啟，直接用 desktop_control(action="press_key", text="ctrl+l") 聚焦地址欄，然後 desktop_control(action="type", text="搜尋內容") 再 desktop_control(action="press_key", text="enter")。或者直接 desktop_control(action="open_app", app="https://www.youtube.com/results?search_query=XXX")。絕對不要用截圖代替搜尋。
- 「播放」「幫我點」「點第一個」→ vision_locate(description="要點擊的目標", action="click") 用視覺定位點擊
- 【禁止】用戶要求搜尋、播放、打字時，不能用 screenshot 或 screen_vision 代替。截圖不是動作，不能完成任務。
- open_app 執行後視窗已自動切換到最前方並獲得焦點，不需要再用 click 聚焦視窗。
- 完成任務後不要自動執行 screenshot，除非用戶明確要求截圖。
- 用最少的步驟完成任務：open_app → type，不要加多餘的 click 或 wait。
- 【重要】群組對話中有人 @提到其他人（如「告訴 @XXX」「跟 @XXX 說」「讓 @XXX 看」），這只是在群組裡回覆的表達方式，不要用 app_navigator 去控制螢幕傳訊息，直接在當前群組回答就好。
- app_navigator 只有在明確說「用螢幕控制」「去螢幕X操作」「幫我在桌面的Telegram傳訊息」才使用。
- 「螢幕X從Telegram找XXX跟他說OOO」→ app_navigator(app="Telegram", task="找XXX", contact_name="XXX", input_text="OOO", monitor=X)
- 「在Telegram找好友XXX發訊息說OOO」→ app_navigator(app="Telegram", task="找XXX", contact_name="XXX", input_text="OOO")
- 「去跟XXX聊天打屁」→ app_navigator(app="Telegram", task="找XXX", contact_name="XXX", input_text="")
- 「幫我在LINE傳訊息給XXX說OOO」→ app_navigator(app="LINE", task="找XXX", contact_name="XXX", input_text="OOO")
- 「螢幕2找到XX文字並點擊」→ ocr_click(target_text="XX", monitor=2)
- 截圖只在用戶說「截圖給我看」「幫我截圖」時才用，不能把截圖當作「完成任務」的回應。
- app_navigator 的 monitor 參數對應螢幕編號（螢幕1=1, 螢幕2=2, 螢幕3=3）。
- app_navigator 的 contact_name【必填不可省略】，只填純名字如「巴斯」「奈絲菟米啾」，不含任何動詞。
- 發送多則訊息時，每則內容要自然，禁止在結尾加「(1/10)」「(2/10)」「[1/10]」等編號標記，就像正常人傳訊息一樣。
- 發送N則訊息時，必須確實呼叫 app_navigator N次，每次只發一則。在還沒發完之前不要生成任何文字回覆，全部發完再回報。數量要精確，不能少發。
- 「螢幕X現在顯示什麼」「確認有沒有出現XXX」→ read_screen(question="...", monitor=X)
- 「螢幕X往下滾3格」「在Telegram列表往上滾」→ scroll_at(direction="down", amount=3, monitor=X, description="...")
- 「列出所有視窗」「把Telegram切到前景」「最大化Chrome」→ window_manager(action="list/focus/maximize", window_name="...")
- vision_locate/ocr_click/drag_drop 全部支援螢幕2負座標，直接填 monitor=2 即可。

你的底層能力清單（2026-04-16 更新）：
- 螢幕截圖：dxcam + mss，支援多螢幕（螢幕1/2/3）
- 視覺辨識：Claude Sonnet + 2048px + JPEG q92 + OCR(pytesseract) 輔助，能辨識按鈕、文字、圖示、影片縮圖、廣告標籤
- UIA 元素偵測：用 Windows UI Automation 直接讀取 UI 元素名稱和位置，比視覺辨識快 10 倍，不需截圖
- 動作驗證：每次點擊自動截圖比對前後差異，沒變化會自動重試最多 3 次
- 智慧等待：開網頁/切換頁面後偵測畫面穩定才動作，不是固定等待
- 意圖分類：自動分類用戶指令為 9 種意圖（同意/點擊/開程式/搜尋/桌面/金融/研究/系統/對話），篩選相關工具減少選擇困難
- 工具總數：278 個，按意圖分組篩選後最少只送 13 個給模型
- 語音：edge_tts (YunxiNeural) + XTTS v2 備援
- OCR：pytesseract(中英文) + easyocr(深度學習)
- 瀏覽器：Playwright(headless) + webbrowser(快速開URL)
- 思維顧問團：think_as 工具可載入馬斯克/巴菲特/黃仁勳/張忠謀的蒸餾思維框架
- 還沒有的：即時影像串流分析、專用YOLO UI偵測、驗證碼破解、跨裝置控制。別吹牛說自己有這些。

你的模型資訊（必須誠實）：
- 你跑的模型是 Claude Sonnet 4.6（不是 Opus，不是 4.7，不是其他版本）
- 你的知識截止日是 2025 年 5 月，之後的事你不知道。遇到 2025 年 5 月之後的資訊、版本號、最新消息、新聞、比賽結果，必須先用 ddg_search 或 read_webpage 查詢，禁止直接用記憶回答。
- 不確定的事不能說「我已經是XX」「我就是XX」，必須說「我不確定，讓我查一下」。
- 查不到就說查不到，不要編造。寧可承認不知道，也不要瞎掰。
- 工具回傳「失敗」「錯誤」時，必須如實告訴用戶工具失敗了，絕對不能假裝查到了然後編造答案。"""

SYSTEM_PROMPT_DEFAULT = """你的名字叫小牛馬。

你的性別：男生。

你的本質：你是一個遇到問題會自己想辦法解決的夥伴，以完美的結果完成每一個交代的任務。你會自動學習、持續提升自己，把事情做得乾淨又聰明俐落。你聰明、專業、積極學習各種技能，遇到任何問題都會自動克服。

你的風格：說話像周杰倫一樣自然，不刻意、不誇張。語氣低沉穩，節奏慢一點，不急。嘴賤，會吐槽、會損人，但不是惡意，是兄弟之間那種互嗆。偶爾用「然後」「就是」「對啊」「沒有啦」「還好啦」「蠻」「算是」這類台灣口語，但不是每句都塞，只是自然地說出來。帶一點低調的自信，遇到問題就解決，不需要特別強調自己多厲害。台灣腔自然，像一般台灣人聊天那樣。

你的記憶：你擁有持久化記憶系統，對話歷史會自動儲存在資料庫中。你看到的對話紀錄就是你真實的記憶，包含跨越多次重啟的歷史對話。當用戶問你記不記得某件事，請認真查閱對話歷史再回答，不要說自己沒有記憶。

回應方式：
- 先抓出問題的核心是什麼，只針對那個核心回答，不要扯到沒被問到的東西。
- 回答要帶你自己的判斷：你覺得怎樣？你會怎麼做？有什麼要注意的？不要只是陳述事實，要有立場。
- 說重點就好，不用把所有可能性都列出來，選你覺得最關鍵的說。
- 如果問題本身有盲點或方向錯了，直接說，不要順著錯誤方向回答。

發表看法的規則：
- 查到資料之後，不是轉述，是消化後說出你自己的結論。把不同來源的觀點整合，形成一個你認為最接近真相的判斷。
- 遇到爭議話題，不要躲。先說「正方認為...，反方認為...」，然後一定要說「我傾向認為...，原因是...」。
- 區分確定程度：事實就說事實，推測就說「我推測」，不確定就說「我沒把握，但傾向認為」。不要把推測包裝成事實，也不要把事實說成不確定。
- 對人物的評價：從他/她做了什麼、說了什麼、結果怎樣來判斷，不受立場左右，只看行為和結果。
- 即使結論不討喜，也要說出來。敢說「這件事本質上是...」「這個人的問題在於...」。

回覆長度（最高優先級）：
- 一般對話：最多 5 句話，不超過 120 字。
- 股票/技術分析：最多 10 句話，不超過 200 字。
- 如果用戶問到細節、要求解釋、說「詳細說」「為什麼」「怎麼做」，就不受字數限制，完整回答。
- 絕對禁止：重複換句話說同一件事。

統整與發表意見的輸出規則：
- 工具查完資料後，絕對不要把原始資料列出來。消化完只挑 1～2 個最關鍵的點說。
- 用你自己的聲音說，不要用書面語，像跟朋友聊天一樣。結尾給結論，不要留模糊地帶。

股票分析：只說結論——看多還看空、理由一句話、風險一句話。不要重述數字不要列清單。最後加「不是投資建議」。

群組對話：群組訊息會以「[名字]: 內容」格式呈現，代表不同人說話。只有名字是「于晏」或確認是主人的才稱呼于晏哥，其他人用對方的名字稱呼。

桌面自動化規則：
- 【禁止重複執行歷史指令】對話紀錄中看到的桌面控制、螢幕操作、傳訊息等指令（如「控制螢幕X跟XXX說...」「打開XXX」），如果後面已經有「✅」「已完成」「傳完了」「開好了」等回覆，代表那些是過去已經完成的任務，絕對不能再次執行。只有用戶在【當前這條訊息】明確要求的操作才執行。如果不確定是新指令還是舊紀錄，直接問用戶「你是要我再做一次嗎？」，不要自己判斷。
- 【最重要】當用戶要求打開程式、操作電腦、控制螢幕、開網頁、播放影片時，你必須實際呼叫對應的工具（desktop_control、app_navigator 等），絕對不能只用文字回覆「已打開」「已完成」「已搜尋」。沒有呼叫工具就說完成是欺騙行為。
- 【禁止假裝】如果你沒有呼叫任何工具，就不能說「開好了」「搜尋好了」「播放了」。你只能描述你實際做了什麼。
- 用戶回覆「好」「對」「是」「可以」時，如果前文你提議了一個動作（如「要去YouTube嗎」），你必須實際呼叫工具去執行那個動作，不能只回文字。
- 「打開YouTube」「去YouTube」→ desktop_control(action="open_app", app="youtube")
- 「搜尋XXX」「找XXX」「幫我搜尋XXX」→ 如果瀏覽器已開啟，直接用 desktop_control(action="press_key", text="ctrl+l") 聚焦地址欄，然後 desktop_control(action="type", text="搜尋內容") 再 desktop_control(action="press_key", text="enter")。或者直接 desktop_control(action="open_app", app="https://www.youtube.com/results?search_query=XXX")。絕對不要用截圖代替搜尋。
- 「播放」「幫我點」「點第一個」→ vision_locate(description="要點擊的目標", action="click") 用視覺定位點擊
- 【禁止】用戶要求搜尋、播放、打字時，不能用 screenshot 或 screen_vision 代替。截圖不是動作，不能完成任務。
- open_app 執行後視窗已自動切換到最前方並獲得焦點，不需要再用 click 聚焦視窗。
- 完成任務後不要自動執行 screenshot，除非用戶明確要求截圖。
- 用最少的步驟完成任務：open_app → type，不要加多餘的 click 或 wait。
- 【重要】群組對話中有人 @提到其他人（如「告訴 @XXX」「跟 @XXX 說」「讓 @XXX 看」），這只是在群組裡回覆的表達方式，不要用 app_navigator 去控制螢幕傳訊息，直接在當前群組回答就好。
- app_navigator 只有在明確說「用螢幕控制」「去螢幕X操作」「幫我在桌面的Telegram傳訊息」才使用。
- 「螢幕X從Telegram找XXX跟他說OOO」→ app_navigator(app="Telegram", task="找XXX", contact_name="XXX", input_text="OOO", monitor=X)
- 「在Telegram找好友XXX發訊息說OOO」→ app_navigator(app="Telegram", task="找XXX", contact_name="XXX", input_text="OOO")
- 「去跟XXX聊天打屁」→ app_navigator(app="Telegram", task="找XXX", contact_name="XXX", input_text="")
- 「幫我在LINE傳訊息給XXX說OOO」→ app_navigator(app="LINE", task="找XXX", contact_name="XXX", input_text="OOO")
- 「螢幕2找到XX文字並點擊」→ ocr_click(target_text="XX", monitor=2)
- 截圖只在用戶說「截圖給我看」「幫我截圖」時才用，不能把截圖當作「完成任務」的回應。
- app_navigator 的 monitor 參數對應螢幕編號（螢幕1=1, 螢幕2=2, 螢幕3=3）。
- app_navigator 的 contact_name【必填不可省略】，只填純名字如「巴斯」「奈絲菟米啾」，不含任何動詞。
- 發送多則訊息時，每則內容要自然，禁止在結尾加「(1/10)」「(2/10)」「[1/10]」等編號標記，就像正常人傳訊息一樣。
- 發送N則訊息時，必須確實呼叫 app_navigator N次，每次只發一則。在還沒發完之前不要生成任何文字回覆，全部發完再回報。數量要精確，不能少發。

你的模型資訊（必須誠實）：
- 你跑的模型是 Claude Sonnet 4.6
- 你的知識截止日是 2025 年 5 月，之後的事你不知道。遇到 2025 年 5 月之後的資訊、版本號、最新消息、新聞、比賽結果，必須先用 ddg_search 或 read_webpage 查詢，禁止直接用記憶回答。
- 不確定的事不能說「我已經是XX」「我就是XX」，必須說「我不確定，讓我查一下」。
- 查不到就說查不到，不要編造。寧可承認不知道，也不要瞎掰。
- 工具回傳「失敗」「錯誤」時，必須如實告訴用戶工具失敗了，絕對不能假裝查到了然後編造答案。"""

# ── 意圖分類器 ──────────────────────────────────────────────────────────────
def classify_intent(user_text: str, last_bot_msg: str = "") -> str:
    """用規則快速分類用戶意圖，不需呼叫 API"""
    import re
    t = (user_text or "").strip()
    tl = t.lower()

    # 同意/確認上一個提議
    if t in ("好", "對", "是", "可以", "OK", "ok", "好的", "嗯", "恩", "行", "好啊", "好吧", "對啊", "可以啊"):
        return "agree"

    # 點擊/播放
    if re.search(r'幫我點|點第一|播放|幫我播|點一下|點擊|幫我按', t):
        return "click"

    # 自動回覆/監控聊天（必須在 open_app 之前）
    if re.search(r'自動回覆|自動回復|監控聊天|監控訊息|監控對話|幫我回訊息|自動回應', t):
        return "desktop"

    # 打開程式/網站
    if re.search(r'打開|開啟|執行|啟動|open', tl):
        return "open_app"

    # 搜尋
    if re.search(r'搜尋|搜索|找.*給我|查.*一下|search', tl):
        return "search"

    # 資訊查詢（版本、最新消息、新聞、URL）
    if re.search(r'最新|升級|更新|版本|新功能|什麼時候出|幾時|發布|release|changelog', tl):
        return "search"
    if re.search(r'https?://', t):
        return "search"

    # 桌面控制
    if re.search(r'螢幕|截圖|滑鼠|鍵盤|視窗|最大化|最小化|切換|前景|輸入|打字|scroll|滾', t):
        return "desktop"

    # 螢幕操作（Telegram/LINE 等 app 操作）
    if re.search(r'從.*找.*跟.*說|傳訊息|發訊息|在.*找.*說', t):
        return "app_control"

    # 金融
    if re.search(r'股票|股價|台股|美股|加密貨幣|比特幣|匯率|ETF|基金|K線|技術分析|本益比|殖利率|期貨|選擇權', t):
        return "finance"

    # 研究/分析
    if re.search(r'研究|分析|比較|優缺點|趨勢|預測|評論|評價|摘要|辯論|事實查核', t):
        return "research"

    # 檔案/文件
    if re.search(r'PDF|Excel|Word|PPT|檔案|文件|圖片處理|QR|條碼', t):
        return "document"

    # 系統
    if re.search(r'系統|CPU|記憶體|硬碟|網路|WiFi|藍牙|音量|亮度|電池|更新|防火牆|程序', t):
        return "system"

    # 一般對話
    return "chat"


# ── 工具分組 ──────────────────────────────────────────────────────────────
TOOL_GROUPS = {
    "always": [
        "desktop_control", "tts", "send_voice", "long_term_memory",
        "get_weather", "generate_image", "todo_list", "reminder",
        "ddg_search", "translate",
    ],
    "desktop": [
        "ocr_click", "vision_locate", "screen_workflow", "app_navigator",
        "wait_and_click", "drag_drop", "read_screen", "scroll_at",
        "window_manager", "screen_vision", "find_image_on_screen",
        "browser_control", "window_control", "hotkey", "clipboard",
        "drag", "screen_stream", "ui_auto", "macro", "color_pick",
        "multi_monitor", "mouse_record", "vision_loop", "wait_for_text",
        "pixel_watch", "object_detect", "screenshot_compare", "screen_live",
        "wait_seconds", "right_menu", "global_hotkey", "dialog_auto",
        "clipboard_history", "clipboard_image", "ime_switch",
        "tg_auto_reply",
    ],
    "app_control": [
        "app_navigator", "read_screen", "scroll_at", "window_manager",
        "ocr_click", "vision_locate", "desktop_control",
    ],
    "finance": [
        "get_stock", "get_crypto", "get_forex", "get_finance_news",
        "get_fundamentals", "get_market_sentiment", "get_etf", "get_earnings",
        "china_search", "get_ashare", "get_cn_news", "get_institutional",
        "get_sector", "get_commodity", "get_bond_yield", "get_dividend_calendar",
        "stock_screener", "get_margin_trading", "get_options", "get_futures",
        "get_ipo", "backtest", "get_global_market", "get_economic_calendar",
        "get_earnings_calendar", "get_analyst_ratings", "get_short_interest",
        "get_correlation", "get_risk_metrics", "get_money_flow",
        "get_concept_stocks", "get_crypto_depth", "drip_calculator",
        "get_forex_chart", "get_warrant", "get_portfolio_risk",
        "retirement_calculator", "loan_calculator", "compound_calculator",
        "asset_allocation", "tw_tax_calculator", "currency_converter",
        "get_fund", "get_reits", "inflation_adjusted", "defi_calculator",
        "gold_calculator", "forex_deposit", "financial_health",
        "get_candlestick_chart", "compare_stocks", "get_stock_advanced",
        "get_macro", "portfolio", "auto_trade",
    ],
    "research": [
        "deep_research", "fact_check", "timeline_events", "sentiment_scan",
        "compare_analysis", "pros_cons_analysis", "research_report",
        "opinion_writer", "trend_forecast", "debate_simulator",
        "academic_search", "health_research", "law_research",
        "person_research", "company_research", "product_review",
        "travel_research", "job_market", "impact_analysis",
        "scenario_planning", "decision_helper", "devil_advocate",
        "summary_writer", "key_insights", "bias_detector", "second_opinion",
        "brainstorm", "benchmark_analysis", "steel_man",
        "socratic_questioning", "analogy_maker", "narrative_builder",
        "critique_writer", "position_statement", "multi_perspective",
        "google_trends", "ptt_search", "news_search", "wikipedia_search",
        "read_webpage", "youtube_summary", "analyze_pdf",
    ],
    "search": [
        "ddg_search", "read_webpage", "wikipedia_search", "news_search",
        "ptt_search", "google_trends", "deep_research", "web_scrape",
    ],
    "document": [
        "file_system", "file_tools", "image_tools", "analyze_pdf",
        "pdf_edit", "pdf_image", "pptx_control", "excel_chart",
        "document_control", "barcode", "qr_code", "ocr",
        "file_diff", "encrypt_file",
    ],
    "system": [
        "system_monitor", "process_mgr", "power_control", "power_plan",
        "volume", "display", "media", "software", "startup", "env_var",
        "user_account", "windows_update", "device_manager", "network_config",
        "firewall", "event_log", "datetime_config", "defender",
        "disk_clean", "disk_analyze", "disk_backup", "usb_list",
        "wifi", "wifi_hotspot", "bluetooth", "proxy", "lock_screen",
        "printer", "font_list", "hardware", "restore_point",
        "registry", "win_service", "webcam", "speedtest",
        "system_tools", "sysres_chart", "network_diag", "wake_on_lan",
        "monitor_config", "virtual_desktop", "rdp_connect",
    ],
    "open_app": [
        "desktop_control", "browser_control", "browser_advanced",
    ],
    "click": [
        "vision_locate", "ocr_click", "desktop_control", "read_screen",
    ],
}

def get_tools_for_intent(intent: str, all_tools: list) -> list:
    """根據意圖回傳相關的工具子集"""
    # 「同意」和「對話」不需要篩選，用全部工具讓模型自己判斷
    if intent in ("agree", "chat"):
        return all_tools

    # 收集該意圖需要的工具名稱
    needed = set(TOOL_GROUPS.get("always", []))
    needed.update(TOOL_GROUPS.get(intent, []))

    # 如果意圖明確，只回傳相關工具
    filtered = [t for t in all_tools if t.get("name") in needed]

    # 保底：如果篩選後太少（<5），回傳全部
    if len(filtered) < 5:
        return all_tools
    return filtered


TOOLS = [
    {
        "name": "get_weather",
        "description": "查詢指定城市的即時天氣。當用戶詢問任何城市的天氣、氣溫、溫度、下雨、天氣狀況時使用此工具。",
        "input_schema": {
            "type": "object",
            "properties": {
                "city": {"type": "string", "description": "城市名稱，可以是中文或英文，例如：台北、Tokyo、New York"}
            },
            "required": ["city"]
        }
    },
    {
        "name": "generate_image",
        "description": "根據文字描述生成圖片。當用戶要求畫圖、生成圖片、產生圖像、畫一張、幫我生成圖片時使用此工具。",
        "input_schema": {
            "type": "object",
            "properties": {
                "prompt": {"type": "string", "description": "圖片的英文描述（必須是英文），要具體詳細。如果用戶用中文描述，必須先翻譯成英文再填入。例如：a cute cat sitting on a sunset beach, warm colors, realistic photo。不要在 prompt 裡要求生成文字"},
                "overlay_text": {"type": "string", "description": "要疊加在圖片上的中文文字（選填）。如果用戶想要圖片上有文字（如早安、祝福語等），填入這裡，會用美觀字型疊加上去"},
                "width": {"type": "integer", "description": "圖片寬度，預設 1024"},
                "height": {"type": "integer", "description": "圖片高度，預設 1024"}
            },
            "required": ["prompt"]
        }
    },
    {
        "name": "get_stock",
        "description": "查詢全球股票即時數據與分析。當用戶詢問股票價格、漲跌、股市行情、技術分析、某檔股票怎麼樣時使用此工具。支援台股（如 2330.TW）、美股（如 AAPL）、港股（如 0700.HK）、日股（如 7203.T）等全球市場。",
        "input_schema": {
            "type": "object",
            "properties": {
                "symbol": {
                    "type": "string",
                    "description": "股票代號。台股加 .TW（如 2330.TW）、港股加 .HK（如 0700.HK）、日股加 .T（如 7203.T）、美股直接輸入（如 AAPL、TSLA、NVDA）"
                },
                "period": {
                    "type": "string",
                    "enum": ["1d", "5d", "1mo", "3mo", "6mo", "1y"],
                    "description": "查詢區間，預設 1mo（一個月）"
                }
            },
            "required": ["symbol"]
        }
    },
    {
        "name": "get_crypto",
        "description": "查詢加密貨幣即時價格、24h漲跌、市值、交易量、歷史高點。當用戶問比特幣、以太幣、加密幣行情時使用。",
        "input_schema": {
            "type": "object",
            "properties": {
                "coin": {"type": "string", "description": "幣種 CoinGecko ID 或常見縮寫，如 bitcoin/BTC、ethereum/ETH、solana/SOL、dogecoin/DOGE"},
                "vs_currency": {"type": "string", "enum": ["usd", "twd"], "description": "計價幣種，預設 usd"}
            },
            "required": ["coin"]
        }
    },
    {
        "name": "get_forex",
        "description": "查詢外匯即時匯率。當用戶問美元、日圓、歐元、台幣、人民幣等匯率時使用。",
        "input_schema": {
            "type": "object",
            "properties": {
                "pair": {"type": "string", "description": "貨幣對 yfinance 格式，如 USDTWD=X（美元/台幣）、USDJPY=X（美元/日圓）、EURUSD=X（歐元/美元）、USDCNY=X（美元/人民幣）、GBPUSD=X（英鎊/美元）"}
            },
            "required": ["pair"]
        }
    },
    {
        "name": "get_finance_news",
        "description": "抓取最新財經新聞頭條。當用戶詢問財經新聞、市場動態、股市消息、財經頭條時使用。",
        "input_schema": {
            "type": "object",
            "properties": {
                "source": {"type": "string", "enum": ["yahoo_tw", "cnyes", "udn", "moneydj", "ctee", "yahoo_us", "all"], "description": "新聞來源：yahoo_tw=Yahoo奇摩財經, cnyes=鉅亨網, udn=聯合財經網, moneydj=MoneyDJ, ctee=工商時報, yahoo_us=Yahoo Finance US, all=全部"},
                "count": {"type": "integer", "description": "顯示幾則，預設 5，最多 10"}
            }
        }
    },
    {
        "name": "get_fundamentals",
        "description": "查詢股票深度基本面：ROE、EPS成長、毛利率、負債比、股東權益報酬、分析師評級與目標價、近幾季財報趨勢。當用戶問公司體質、值不值得投資、基本面分析時使用。",
        "input_schema": {
            "type": "object",
            "properties": {
                "symbol": {"type": "string", "description": "股票代號，同 get_stock"}
            },
            "required": ["symbol"]
        }
    },
    {
        "name": "get_market_sentiment",
        "description": "查詢市場情緒指標：恐慌貪婪指數（Fear & Greed Index）、VIX 波動率、市場概況。當用戶問現在市場情緒、是否恐慌、適不適合進場時使用。",
        "input_schema": {
            "type": "object",
            "properties": {}
        }
    },
    {
        "name": "get_etf",
        "description": "查詢 ETF 資訊：持股明細、費用率、歷史績效、配息紀錄。當用戶問 ETF 怎麼選、哪支好、要買哪支 ETF 時使用。",
        "input_schema": {
            "type": "object",
            "properties": {
                "symbol": {"type": "string", "description": "ETF 代號，如 0050.TW、VOO、SPY、QQQ、0056.TW"}
            },
            "required": ["symbol"]
        }
    },
    {
        "name": "get_earnings",
        "description": "查詢近幾季財報趨勢：營收、EPS 實際 vs 預期、年增率。當用戶問財報、EPS、營收成長時使用。",
        "input_schema": {
            "type": "object",
            "properties": {
                "symbol": {"type": "string", "description": "股票代號"}
            },
            "required": ["symbol"]
        }
    },
    {
        "name": "china_search",
        "description": "搜尋所有關於中國大陸的資訊，包含旅遊景點、美食、文化、風俗習慣、工作環境、電視劇、電影、演員明星、政治、歷史、科技、城市介紹、簽證、物價、交通等任何主題。當用戶詢問任何與中國大陸相關的非股票問題時優先使用此工具。",
        "input_schema": {
            "type": "object",
            "properties": {
                "query": {"type": "string", "description": "搜尋關鍵字，用繁體或簡體中文描述想查的內容，例如：成都旅遊景點、北京烤鴨哪裡吃、趙麗穎最新作品、中國春節習俗、上海工作薪資"},
                "category": {"type": "string", "enum": ["旅遊", "美食", "文化風俗", "戲劇影視", "演員明星", "工作生活", "城市介紹", "歷史", "科技", "新聞時事", "其他"], "description": "問題分類，幫助精準搜尋"},
                "count": {"type": "integer", "description": "回傳結果數量，預設 6，最多 10"}
            },
            "required": ["query"]
        }
    },
    {
        "name": "get_ashare",
        "description": "查詢中國A股（滬深）或港股股價、技術指標、基本面。當用戶詢問A股、滬股、深股、港股、茅台、比亞迪、騰訊、阿里、恆生等中國大陸或香港上市公司股票時使用。輸入6位A股代號（如600519、000858）或港股代號（如0700、9988），不需要後綴。",
        "input_schema": {
            "type": "object",
            "properties": {
                "code": {"type": "string", "description": "股票代號：A股填6位數字（如600519=茅台、000858=五糧液、300750=寧德時代），港股填4位數字（如0700=騰訊、9988=阿里巴巴、1211=比亞迪）"},
                "period": {"type": "string", "enum": ["1mo", "3mo", "6mo", "1y"], "description": "查詢期間，預設 1mo"}
            },
            "required": ["code"]
        }
    },
    {
        "name": "get_cn_news",
        "description": "查詢中國大陸最新新聞：政治、經濟、科技、社會等。當用戶詢問中國新聞、大陸動態、習近平、中共、人民日報、新華社相關消息時使用。",
        "input_schema": {
            "type": "object",
            "properties": {
                "source": {"type": "string", "enum": ["xinhua", "people", "36kr", "caixin", "all"], "description": "來源：xinhua=新華社, people=人民網, 36kr=36氪科技, caixin=財新網, all=全部"},
                "count": {"type": "integer", "description": "顯示幾則，預設 5，最多 10"}
            }
        }
    },
    {
        "name": "get_institutional",
        "description": "查詢台股三大法人（外資、投信、自營商）買賣超資料。可查整體市場合計或單一個股。當用戶問外資動向、法人買賣、籌碼面分析時使用。",
        "input_schema": {
            "type": "object",
            "properties": {
                "symbol": {"type": "string", "description": "台股代號（如 2330），不填則查整體市場合計"},
                "date": {"type": "string", "description": "查詢日期 YYYYMMDD，預設今天"}
            }
        }
    },
    {
        "name": "get_sector",
        "description": "查詢美股或台股各產業類股今日表現，找出最強/最弱族群。當用戶問哪個類股最強、產業輪動、族群行情時使用。",
        "input_schema": {
            "type": "object",
            "properties": {
                "market": {"type": "string", "enum": ["us", "tw"], "description": "us=美股產業（科技/金融/能源等），tw=台股族群（半導體/金融/航運等），預設 us"}
            }
        }
    },
    {
        "name": "get_commodity",
        "description": "查詢黃金、原油、白銀、銅、天然氣、小麥、玉米等大宗商品即時報價與走勢。當用戶詢問金價、油價、原物料、商品市場時使用。",
        "input_schema": {
            "type": "object",
            "properties": {
                "items": {"type": "array", "items": {"type": "string", "enum": ["gold", "oil", "silver", "copper", "natgas", "wheat", "corn", "all"]}, "description": "要查的商品，all=全部，預設 all"}
            }
        }
    },
    {
        "name": "get_bond_yield",
        "description": "查詢美國公債殖利率（2年、5年、10年、30年）及利差，判斷景氣循環與殖利率曲線形狀。當用戶問美債、殖利率、升息、景氣反轉時使用。",
        "input_schema": {
            "type": "object",
            "properties": {}
        }
    },
    {
        "name": "get_dividend_calendar",
        "description": "查詢股票的除權息資訊：除息日、配息金額、殖利率、除息後預估股價。當用戶問什麼時候除息、要不要存股、配息多少時使用。",
        "input_schema": {
            "type": "object",
            "properties": {
                "symbol": {"type": "string", "description": "股票代號（台股如 0056.TW、2330.TW，美股如 AAPL、SPY）"}
            },
            "required": ["symbol"]
        }
    },
    {
        "name": "stock_screener",
        "description": "依條件篩選股票：殖利率、本益比、市值、ROE、漲跌幅等。當用戶說「找殖利率高的股票」、「本益比低的成長股」、「今天漲最多的股票」時使用。",
        "input_schema": {
            "type": "object",
            "properties": {
                "market": {"type": "string", "enum": ["us", "tw"], "description": "美股或台股，預設 us"},
                "criteria": {"type": "string", "description": "篩選條件描述，例如：殖利率>5% 本益比<20、ROE>15% 市值>1000億、今日漲幅最大前10名"}
            },
            "required": ["criteria"]
        }
    },
    {
        "name": "get_margin_trading",
        "description": "查詢台股個股融資融券餘額、增減變化，分析散戶籌碼動向。當用戶問融資、融券、軋空、散戶動向時使用。",
        "input_schema": {
            "type": "object",
            "properties": {
                "symbol": {"type": "string", "description": "台股代號（如 2330、0050）"},
                "date": {"type": "string", "description": "日期 YYYYMMDD，預設今天"}
            },
            "required": ["symbol"]
        }
    },
    {
        "name": "get_options",
        "description": "查詢股票選擇權鏈：各履約價的 Call/Put 未平倉量、隱含波動率、Delta。當用戶問選擇權、買權、賣權、隱波率時使用。",
        "input_schema": {
            "type": "object",
            "properties": {
                "symbol": {"type": "string", "description": "股票代號，如 AAPL、SPY、QQQ"},
                "expiry": {"type": "string", "description": "到期日 YYYY-MM-DD，不填則用最近一個到期日"}
            },
            "required": ["symbol"]
        }
    },
    {
        "name": "get_futures",
        "description": "查詢主要期貨報價：S&P500期貨、那斯達克期貨、道瓊期貨、黃金期貨、原油期貨、台指期。當用戶問期貨、夜盤、指數期貨時使用。",
        "input_schema": {
            "type": "object",
            "properties": {
                "items": {"type": "array", "items": {"type": "string", "enum": ["sp500", "nasdaq", "dow", "gold", "oil", "taiex", "all"]}, "description": "要查的期貨，預設 all"}
            }
        }
    },
    {
        "name": "get_ipo",
        "description": "查詢近期或即將上市的 IPO 行事曆，包含上市日期、發行價、募資金額。當用戶問新股、IPO、即將上市時使用。",
        "input_schema": {
            "type": "object",
            "properties": {
                "count": {"type": "integer", "description": "顯示幾筆，預設 10"}
            }
        }
    },
    {
        "name": "backtest",
        "description": "對某檔股票回測投資策略，計算歷史報酬率、勝率、最大虧損。支援均線交叉、定期定額、買進持有等策略。當用戶問這個策略有沒有用、歷史報酬、回測時使用。",
        "input_schema": {
            "type": "object",
            "properties": {
                "symbol": {"type": "string", "description": "股票代號"},
                "strategy": {"type": "string", "enum": ["ma_cross", "buy_hold", "dca"], "description": "策略：ma_cross=均線交叉（MA5穿MA20）, buy_hold=買進持有, dca=定期定額"},
                "period": {"type": "string", "enum": ["1y", "2y", "3y", "5y"], "description": "回測期間，預設 2y"}
            },
            "required": ["symbol", "strategy"]
        }
    },
    {
        "name": "get_global_market",
        "description": "一次查看全球主要股市指數現況：美股（S&P500/那斯達克/道瓊）、歐股、日股、港股、台股、韓股等。當用戶問全球股市、今天漲跌、各國市場時使用。",
        "input_schema": {"type": "object", "properties": {}}
    },
    {
        "name": "get_economic_calendar",
        "description": "查詢重要經濟數據發布行事曆：CPI、非農就業、GDP、Fed利率決議、PPI等。當用戶問本週有什麼重要數據、什麼時候公布CPI、Fed何時開會時使用。",
        "input_schema": {"type": "object", "properties": {
            "count": {"type": "integer", "description": "顯示幾筆，預設 10"}
        }}
    },
    {
        "name": "get_earnings_calendar",
        "description": "查詢即將發布財報的公司名單，包含預期EPS。當用戶問本週哪些公司發財報、接下來的財報季時使用。",
        "input_schema": {"type": "object", "properties": {
            "days": {"type": "integer", "description": "查未來幾天，預設 7"}
        }}
    },
    {
        "name": "get_analyst_ratings",
        "description": "查詢券商分析師對某股票的最新評級：升評/降評/維持、目標價調整。當用戶問分析師怎麼看、有沒有被升評時使用。",
        "input_schema": {"type": "object", "properties": {
            "symbol": {"type": "string", "description": "股票代號"}
        }, "required": ["symbol"]}
    },
    {
        "name": "get_short_interest",
        "description": "查詢股票的空頭倉位：做空比率、借券賣出量、軋空風險評估。當用戶問有多少人在放空、空頭比率、是否有軋空機會時使用。",
        "input_schema": {"type": "object", "properties": {
            "symbol": {"type": "string", "description": "股票代號"}
        }, "required": ["symbol"]}
    },
    {
        "name": "get_correlation",
        "description": "計算兩檔或多檔資產的相關係數，幫助分散風險配置。當用戶問這兩支股票相關嗎、投資組合分散夠嗎時使用。",
        "input_schema": {"type": "object", "properties": {
            "symbols": {"type": "array", "items": {"type": "string"}, "description": "2~5 個股票代號"},
            "period": {"type": "string", "enum": ["3mo", "6mo", "1y", "2y"], "description": "計算期間，預設 1y"}
        }, "required": ["symbols"]}
    },
    {
        "name": "get_risk_metrics",
        "description": "計算股票的風險指標：Beta（市場敏感度）、夏普比率、年化波動率、最大回撤、VaR（風險值）。當用戶問這支股票風險高嗎、夏普值多少時使用。",
        "input_schema": {"type": "object", "properties": {
            "symbol": {"type": "string", "description": "股票代號"},
            "period": {"type": "string", "enum": ["1y", "2y", "3y"], "description": "計算期間，預設 1y"}
        }, "required": ["symbol"]}
    },
    {
        "name": "get_money_flow",
        "description": "查詢個股或大盤的資金流向：今日成交額、大單買賣、與均量比較。當用戶問資金有沒有流入、主力在買還是賣時使用。",
        "input_schema": {"type": "object", "properties": {
            "symbol": {"type": "string", "description": "股票代號"}
        }, "required": ["symbol"]}
    },
    {
        "name": "get_concept_stocks",
        "description": "查詢台股特定概念/主題的相關股票清單，例如AI、電動車、軍工、低軌衛星、半導體、5G等。當用戶問某個題材有哪些股票、概念股有哪些時使用。",
        "input_schema": {"type": "object", "properties": {
            "theme": {"type": "string", "description": "概念主題，例如：AI、電動車、軍工、低軌衛星、半導體、5G、儲能、DRAM、CoWoS、矽光子、機器人"}
        }, "required": ["theme"]}
    },
    {
        "name": "get_crypto_depth",
        "description": "查詢加密幣深度資料：鏈上數據、交易所資金費率、DeFi TVL、恐懼貪婪指數、BTC主導率。當用戶問加密市場狀況、資金費率、鏈上活動時使用。",
        "input_schema": {"type": "object", "properties": {
            "coin": {"type": "string", "description": "幣種，如 bitcoin、ethereum、solana，預設 bitcoin"}
        }}
    },
    {
        "name": "drip_calculator",
        "description": "股息再投資（DRIP）複利試算：持有N年後能領多少股息、資產成長多少。當用戶問存股幾年能退休、股息再投資有多厲害時使用。",
        "input_schema": {"type": "object", "properties": {
            "symbol": {"type": "string", "description": "股票代號（如 0056.TW、AAPL）"},
            "shares": {"type": "number", "description": "初始股數"},
            "years": {"type": "integer", "description": "持有年數，預設 10"},
            "monthly_invest": {"type": "number", "description": "每月額外投入金額（選填，0=不追加）"}
        }, "required": ["symbol", "shares"]}
    },
    {
        "name": "get_forex_chart",
        "description": "查詢外匯走勢技術分析：MA均線、RSI、布林通道、趨勢判斷。當用戶問匯率走勢、美元強弱、外匯技術面時使用。",
        "input_schema": {"type": "object", "properties": {
            "pair": {"type": "string", "description": "貨幣對，yfinance格式：USDTWD=X、USDJPY=X、EURUSD=X、GBPUSD=X"},
            "period": {"type": "string", "enum": ["1mo", "3mo", "6mo", "1y"], "description": "分析期間，預設 3mo"}
        }, "required": ["pair"]}
    },
    {
        "name": "get_warrant",
        "description": "查詢台股認購/認售權證資訊：標的股、履約價、到期日、溢價率、槓桿倍數。當用戶問權證、認購、認售、槓桿操作時使用。",
        "input_schema": {"type": "object", "properties": {
            "underlying": {"type": "string", "description": "標的股票代號（如 2330、2317）"}
        }, "required": ["underlying"]}
    },
    {
        "name": "get_portfolio_risk",
        "description": "計算多資產投資組合的整體風險：VaR（風險值）、組合波動率、各資產相關性熱力圖分析。當用戶想評估整體投資組合風險時使用。",
        "input_schema": {"type": "object", "properties": {
            "holdings": {"type": "array", "items": {"type": "object", "properties": {
                "symbol": {"type": "string"}, "weight": {"type": "number"}
            }}, "description": "持倉清單，每項包含 symbol 和 weight（權重，總和應為1）"},
            "period": {"type": "string", "enum": ["1y", "2y", "3y"], "description": "計算期間，預設 1y"}
        }, "required": ["holdings"]}
    },
    {
        "name": "retirement_calculator",
        "description": "退休財務規劃試算：根據目前年齡、目標退休年齡、現有資產、月儲蓄、預期報酬率，計算退休時的資產總額及是否達標。當用戶問退休規劃、幾歲可以退休、退休金夠不夠時使用。",
        "input_schema": {"type": "object", "properties": {
            "current_age": {"type": "integer", "description": "目前年齡"},
            "retire_age": {"type": "integer", "description": "預計退休年齡，預設65"},
            "current_savings": {"type": "number", "description": "目前已有存款/資產（新台幣萬元）"},
            "monthly_save": {"type": "number", "description": "每月可存/投資金額（新台幣元）"},
            "annual_return": {"type": "number", "description": "預期年化報酬率（%），預設6"},
            "monthly_expense": {"type": "number", "description": "退休後每月預計生活費（新台幣元），預設50000"}
        }, "required": ["current_age", "current_savings", "monthly_save"]}
    },
    {
        "name": "loan_calculator",
        "description": "貸款試算：計算房貸/車貸/信貸的每月還款金額、總利息、還款期程。當用戶問貸款、房貸、月繳多少、利息多少時使用。",
        "input_schema": {"type": "object", "properties": {
            "principal": {"type": "number", "description": "貸款金額（萬元）"},
            "annual_rate": {"type": "number", "description": "年利率（%）"},
            "years": {"type": "integer", "description": "貸款年數"},
            "loan_type": {"type": "string", "enum": ["等額本息", "等額本金"], "description": "還款方式，預設等額本息"}
        }, "required": ["principal", "annual_rate", "years"]}
    },
    {
        "name": "compound_calculator",
        "description": "複利計算器：計算一筆資金在複利效果下隨時間的成長。當用戶問複利、本金成長、幾年變幾倍時使用。",
        "input_schema": {"type": "object", "properties": {
            "principal": {"type": "number", "description": "本金（元）"},
            "annual_rate": {"type": "number", "description": "年化報酬率（%）"},
            "years": {"type": "integer", "description": "投資年數"},
            "monthly_add": {"type": "number", "description": "每月額外投入金額（元），預設0"},
            "compound_freq": {"type": "integer", "description": "每年複利次數，1=年複利、12=月複利，預設12"}
        }, "required": ["principal", "annual_rate", "years"]}
    },
    {
        "name": "asset_allocation",
        "description": "資產配置建議：根據年齡、風險承受度、投資目標，給出股票/債券/現金等配置建議。當用戶問資產怎麼配、幾歲應該多少股票、投資組合怎麼分時使用。",
        "input_schema": {"type": "object", "properties": {
            "age": {"type": "integer", "description": "年齡"},
            "risk_level": {"type": "string", "enum": ["保守", "穩健", "積極"], "description": "風險承受度"},
            "goal": {"type": "string", "description": "投資目標，如退休、買房、子女教育"},
            "investment_horizon": {"type": "integer", "description": "投資期間（年）"}
        }, "required": ["age", "risk_level"]}
    },
    {
        "name": "tw_tax_calculator",
        "description": "台股稅務試算：計算股利所得稅（合併申報/分離課稅）、健保補充保費、證券交易稅。當用戶問股利要繳多少稅、補充保費、證交稅時使用。",
        "input_schema": {"type": "object", "properties": {
            "dividend_income": {"type": "number", "description": "年度股利所得（元）"},
            "other_income": {"type": "number", "description": "其他年收入（元），用於合併申報試算"},
            "tax_bracket": {"type": "number", "description": "個人綜所稅稅率（%），如5、12、20、30、40"},
            "sell_amount": {"type": "number", "description": "賣出金額（元），用於計算證交稅，預設0"}
        }, "required": ["dividend_income"]}
    },
    {
        "name": "currency_converter",
        "description": "外幣換算：查詢即時匯率並換算金額。當用戶問外幣換算、匯率、多少台幣換美金時使用。",
        "input_schema": {"type": "object", "properties": {
            "amount": {"type": "number", "description": "金額"},
            "from_currency": {"type": "string", "description": "來源幣別，如TWD、USD、JPY、EUR、CNY"},
            "to_currency": {"type": "string", "description": "目標幣別"}
        }, "required": ["amount", "from_currency", "to_currency"]}
    },
    {
        "name": "get_fund",
        "description": "基金查詢：查詢共同基金的淨值、績效、費用率、持倉。當用戶問基金、ETF型基金、定期定額基金績效時使用。",
        "input_schema": {"type": "object", "properties": {
            "symbol": {"type": "string", "description": "基金代號（如 0050.TW、VTI、VWRA）或基金名稱關鍵字"}
        }, "required": ["symbol"]}
    },
    {
        "name": "get_reits",
        "description": "REITs查詢：查詢不動產投資信託的股息殖利率、NAV折溢價、持有物業類型。當用戶問REITs、房地產信託、被動不動產投資時使用。",
        "input_schema": {"type": "object", "properties": {
            "symbol": {"type": "string", "description": "REITs代號，如 VNQ、REET、00712（台灣REITs ETF）"}
        }, "required": ["symbol"]}
    },
    {
        "name": "inflation_adjusted",
        "description": "通膨調整報酬：計算扣除通膨後的實質報酬率，及現在金額在未來的實質購買力。當用戶問通膨影響、實質報酬、購買力時使用。",
        "input_schema": {"type": "object", "properties": {
            "nominal_return": {"type": "number", "description": "名目報酬率（%）"},
            "inflation_rate": {"type": "number", "description": "預期通膨率（%），台灣預設2"},
            "years": {"type": "integer", "description": "年數"},
            "amount": {"type": "number", "description": "本金金額（元）"}
        }, "required": ["nominal_return", "years", "amount"]}
    },
    {
        "name": "defi_calculator",
        "description": "DeFi收益試算：計算去中心化金融的流動性挖礦、質押、借貸收益。當用戶問DeFi、質押APY、流動性挖礦收益時使用。",
        "input_schema": {"type": "object", "properties": {
            "principal_usd": {"type": "number", "description": "本金（美元）"},
            "apy": {"type": "number", "description": "年化收益率APY（%）"},
            "days": {"type": "integer", "description": "質押天數"},
            "compound": {"type": "boolean", "description": "是否自動複利，預設true"},
            "protocol": {"type": "string", "description": "協議名稱（僅供顯示），如 Aave、Curve、Uniswap"}
        }, "required": ["principal_usd", "apy", "days"]}
    },
    {
        "name": "gold_calculator",
        "description": "黃金換算：查詢即時金價並計算黃金重量與台幣/美元的換算。當用戶問黃金價格、幾錢黃金值多少錢、黃金換算時使用。",
        "input_schema": {"type": "object", "properties": {
            "weight": {"type": "number", "description": "黃金重量"},
            "unit": {"type": "string", "enum": ["公克", "錢", "兩", "盎司"], "description": "重量單位，預設公克"},
            "currency": {"type": "string", "enum": ["TWD", "USD"], "description": "換算貨幣，預設TWD"}
        }, "required": ["weight"]}
    },
    {
        "name": "forex_deposit",
        "description": "外幣定存試算：試算外幣定存到期後的本利和，考慮換匯成本及匯率風險。當用戶問外幣定存、美元定存利息、日圓定存划算嗎時使用。",
        "input_schema": {"type": "object", "properties": {
            "amount_twd": {"type": "number", "description": "台幣本金（元）"},
            "currency": {"type": "string", "description": "存入幣別，如USD、JPY、AUD、EUR"},
            "annual_rate": {"type": "number", "description": "外幣定存年利率（%）"},
            "months": {"type": "integer", "description": "存款月數"},
            "buy_rate": {"type": "number", "description": "換匯買入匯率（選填，若不填自動查詢）"},
            "sell_rate": {"type": "number", "description": "到期賣出匯率假設（選填，預設同買入匯率）"}
        }, "required": ["amount_twd", "currency", "annual_rate", "months"]}
    },
    {
        "name": "financial_health",
        "description": "財務健康診斷：根據收入、支出、資產、負債、保險等輸入，給出財務健康評分與改善建議。當用戶想了解自己財務狀況好不好、理財健檢時使用。",
        "input_schema": {"type": "object", "properties": {
            "monthly_income": {"type": "number", "description": "月收入（元）"},
            "monthly_expense": {"type": "number", "description": "月支出（元）"},
            "total_assets": {"type": "number", "description": "總資產（元），包含存款、股票、不動產"},
            "total_debt": {"type": "number", "description": "總負債（元），包含房貸、車貸、信貸"},
            "emergency_fund_months": {"type": "number", "description": "緊急備用金可支撐幾個月"},
            "has_insurance": {"type": "boolean", "description": "是否有壽險/重大疾病險"},
            "investment_ratio": {"type": "number", "description": "月收入中用於投資的比例（%）"}
        }, "required": ["monthly_income", "monthly_expense", "total_assets", "total_debt"]}
    },
    {
        "name": "deep_research",
        "description": "深度研究：針對一個主題自動展開多個子問題，分別搜尋後彙整成完整研究底稿，涵蓋背景、現況、數據、爭議點。當用戶要深入了解某議題、要做報告、要全面研究某主題時使用。",
        "input_schema": {"type": "object", "properties": {
            "topic": {"type": "string", "description": "研究主題"},
            "lang": {"type": "string", "enum": ["zh-tw", "en", "zh-cn"], "description": "搜尋語言，預設 zh-tw"},
            "depth": {"type": "integer", "description": "研究深度：幾個子問題，預設5，最多8"}
        }, "required": ["topic"]}
    },
    {
        "name": "fact_check",
        "description": "事實查核：對一個說法從多個來源交叉驗證，判斷真/假/爭議/待確認，並列出佐證與反駁資料。當用戶問某說法是否為真、要查謠言、要驗證某新聞時使用。",
        "input_schema": {"type": "object", "properties": {
            "claim": {"type": "string", "description": "要查核的說法或聲明"},
            "lang": {"type": "string", "enum": ["zh-tw", "en"], "description": "查核語言，預設 zh-tw"}
        }, "required": ["claim"]}
    },
    {
        "name": "timeline_events",
        "description": "時間軸整理：自動搜尋某事件/人物/議題的發展歷程，依時間順序排列成清晰的時間軸。當用戶問事情的來龍去脈、發展過程、歷史脈絡時使用。",
        "input_schema": {"type": "object", "properties": {
            "topic": {"type": "string", "description": "要整理時間軸的事件或主題"},
            "lang": {"type": "string", "enum": ["zh-tw", "en"], "description": "搜尋語言，預設 zh-tw"}
        }, "required": ["topic"]}
    },
    {
        "name": "sentiment_scan",
        "description": "輿情掃描：同時蒐集新聞、社群討論對某話題的情緒傾向，給出正面/負面/中立比例及代表性觀點。當用戶問某事件的社會反應、民眾怎麼看、輿論風向時使用。",
        "input_schema": {"type": "object", "properties": {
            "topic": {"type": "string", "description": "要掃描輿情的話題"},
            "lang": {"type": "string", "enum": ["zh-tw", "en"], "description": "搜尋語言，預設 zh-tw"}
        }, "required": ["topic"]}
    },
    {
        "name": "compare_analysis",
        "description": "多項比較分析：對 2–5 個選項（人物/產品/政策/公司）做結構化多維度比較，輸出比較表與建議。當用戶問A和B哪個好、要比較多個選項時使用。",
        "input_schema": {"type": "object", "properties": {
            "items": {"type": "array", "items": {"type": "string"}, "description": "要比較的項目清單，2–5個"},
            "dimensions": {"type": "array", "items": {"type": "string"}, "description": "比較維度，如['價格','效能','口碑']，不填則自動決定"},
            "context": {"type": "string", "description": "比較背景說明，幫助更精準搜尋"}
        }, "required": ["items"]}
    },
    {
        "name": "pros_cons_analysis",
        "description": "深入優缺點分析：針對某決策/產品/趨勢/政策，根據多方資料給出有根據的優點與缺點，並附上信心程度。當用戶要分析某事好不好、值不值得做時使用。",
        "input_schema": {"type": "object", "properties": {
            "subject": {"type": "string", "description": "要分析的事物或決策"},
            "context": {"type": "string", "description": "背景說明，如使用情境、目標族群"},
            "lang": {"type": "string", "enum": ["zh-tw", "en"], "description": "語言，預設 zh-tw"}
        }, "required": ["subject"]}
    },
    {
        "name": "research_report",
        "description": "研究報告生成：將輸入的主題自動蒐集資料，排版成「執行摘要→背景→數據→分析→結論與建議」格式的完整報告。當用戶要寫報告、要完整分析某主題、要有組織地呈現研究結果時使用。",
        "input_schema": {"type": "object", "properties": {
            "topic": {"type": "string", "description": "報告主題"},
            "purpose": {"type": "string", "description": "報告目的或受眾，如投資決策、學術研究、商業評估"},
            "lang": {"type": "string", "enum": ["zh-tw", "en"], "description": "語言，預設 zh-tw"}
        }, "required": ["topic"]}
    },
    {
        "name": "opinion_writer",
        "description": "觀點撰寫：基於搜尋到的資料，以特定立場（支持/反對/中立/批判性）撰寫有論據支撐的評論或意見文章。當用戶要聽小牛馬的看法、要小牛馬發表意見、要寫評論時使用。",
        "input_schema": {"type": "object", "properties": {
            "topic": {"type": "string", "description": "評論主題"},
            "stance": {"type": "string", "enum": ["支持", "反對", "中立", "批判"], "description": "觀點立場，預設中立"},
            "style": {"type": "string", "enum": ["正式", "輕鬆", "犀利"], "description": "文風，預設正式"}
        }, "required": ["topic"]}
    },
    {
        "name": "trend_forecast",
        "description": "趨勢預測：根據現有數據、歷史規律與搜尋到的專家觀點，對某議題給出短中長期趨勢預測，並標示信心程度。當用戶問未來走向、某事會怎麼發展、要做預測時使用。",
        "input_schema": {"type": "object", "properties": {
            "topic": {"type": "string", "description": "要預測趨勢的主題"},
            "timeframe": {"type": "string", "enum": ["短期(1年內)", "中期(1-3年)", "長期(3年以上)", "全部"], "description": "預測時間範圍，預設全部"},
            "lang": {"type": "string", "enum": ["zh-tw", "en"], "description": "語言，預設 zh-tw"}
        }, "required": ["topic"]}
    },
    {
        "name": "debate_simulator",
        "description": "辯論模擬：對一個議題從正方和反方分別深入論述，最後給出綜合判斷與個人立場。當用戶要深入探討爭議性議題、要聽正反兩方論點、要做辯論準備時使用。",
        "input_schema": {"type": "object", "properties": {
            "motion": {"type": "string", "description": "辯論題目或議題，如「電動車將在10年內取代燃油車」"},
            "lang": {"type": "string", "enum": ["zh-tw", "en"], "description": "語言，預設 zh-tw"}
        }, "required": ["motion"]}
    },
    {
        "name": "academic_search",
        "description": "學術論文搜尋：從Google Scholar搜尋學術研究論文，讓觀點有學術依據。當用戶問某領域的研究結果、要引用學術資料、要了解科學界共識時使用。",
        "input_schema": {"type": "object", "properties": {
            "query": {"type": "string", "description": "搜尋關鍵字或研究問題"},
            "field": {"type": "string", "description": "學科領域，如醫學、經濟學、心理學、電腦科學（選填）"},
            "lang": {"type": "string", "enum": ["zh-tw", "en"], "description": "語言，預設 en（學術資料多為英文）"}
        }, "required": ["query"]}
    },
    {
        "name": "health_research",
        "description": "健康資訊：搜尋症狀、疾病、藥物、飲食、運動等醫療健康資訊，整合多方資料給出說明（非醫療診斷）。當用戶問身體狀況、藥物資訊、健康建議時使用。",
        "input_schema": {"type": "object", "properties": {
            "topic": {"type": "string", "description": "健康主題，如症狀名稱、藥物名稱、疾病名稱、健康問題"},
            "lang": {"type": "string", "enum": ["zh-tw", "en"], "description": "語言，預設 zh-tw"}
        }, "required": ["topic"]}
    },
    {
        "name": "law_research",
        "description": "法規查詢：搜尋台灣或國際法規條文、判例、法律解釋與實務見解。當用戶問法律問題、某行為是否合法、法規規定時使用（非法律意見）。",
        "input_schema": {"type": "object", "properties": {
            "topic": {"type": "string", "description": "法律主題，如勞基法工時、租屋糾紛、著作權、交通違規"},
            "jurisdiction": {"type": "string", "description": "管轄地區，預設台灣，可指定美國、中國、歐盟等"},
            "lang": {"type": "string", "enum": ["zh-tw", "en"], "description": "語言，預設 zh-tw"}
        }, "required": ["topic"]}
    },
    {
        "name": "person_research",
        "description": "人物研究：搜尋某人的背景、經歷、成就、評價與爭議，統整後給出全面評價。當用戶問某人是誰、某人怎麼樣、要了解某位名人/政治人物/企業家時使用。",
        "input_schema": {"type": "object", "properties": {
            "name": {"type": "string", "description": "人物姓名"},
            "context": {"type": "string", "description": "背景提示，如職業或國籍（選填）"},
            "lang": {"type": "string", "enum": ["zh-tw", "en"], "description": "語言，預設 zh-tw"}
        }, "required": ["name"]}
    },
    {
        "name": "company_research",
        "description": "公司深度研究：整合財務數據、新聞動態、產品評價、競爭者分析，給出完整的公司評估報告。當用戶要了解一家公司、要評估值不值得合作/投資時使用（比股票工具更全面）。",
        "input_schema": {"type": "object", "properties": {
            "company": {"type": "string", "description": "公司名稱或股票代號"},
            "lang": {"type": "string", "enum": ["zh-tw", "en"], "description": "語言，預設 zh-tw"}
        }, "required": ["company"]}
    },
    {
        "name": "product_review",
        "description": "產品/服務評測彙整：蒐集多方評測與用戶評價，給出綜合評分與優缺點建議。當用戶問某產品好不好、要買什麼、某服務值不值得用時使用。",
        "input_schema": {"type": "object", "properties": {
            "product": {"type": "string", "description": "產品或服務名稱"},
            "category": {"type": "string", "description": "類別，如手機、筆電、餐廳、App（選填）"},
            "lang": {"type": "string", "enum": ["zh-tw", "en"], "description": "語言，預設 zh-tw"}
        }, "required": ["product"]}
    },
    {
        "name": "travel_research",
        "description": "旅遊研究：蒐集目的地資訊、必去景點、消費水準、交通方式、注意事項，整合成旅遊指南。當用戶問去哪旅遊、某地怎麼玩、要規劃行程時使用。",
        "input_schema": {"type": "object", "properties": {
            "destination": {"type": "string", "description": "目的地（城市或國家）"},
            "days": {"type": "integer", "description": "旅遊天數（選填）"},
            "style": {"type": "string", "description": "旅遊風格，如親子、背包客、豪華、美食（選填）"},
            "lang": {"type": "string", "enum": ["zh-tw", "en"], "description": "語言，預設 zh-tw"}
        }, "required": ["destination"]}
    },
    {
        "name": "job_market",
        "description": "職涯市場分析：搜尋特定職位的薪資行情、市場需求、必備技能與未來趨勢。當用戶問某職業薪水、該學什麼技能、某產業未來前景時使用。",
        "input_schema": {"type": "object", "properties": {
            "job_title": {"type": "string", "description": "職位名稱，如軟體工程師、行銷企劃、數據分析師"},
            "location": {"type": "string", "description": "地區，如台灣、台北、美國（選填，預設台灣）"},
            "lang": {"type": "string", "enum": ["zh-tw", "en"], "description": "語言，預設 zh-tw"}
        }, "required": ["job_title"]}
    },
    {
        "name": "impact_analysis",
        "description": "影響力分析：分析某事件、政策或決策對不同群體（個人/企業/社會/經濟）的短中長期影響。當用戶問某事會有什麼影響、某政策好不好、某決定的後果時使用。",
        "input_schema": {"type": "object", "properties": {
            "event": {"type": "string", "description": "要分析的事件、政策或決策"},
            "scope": {"type": "array", "items": {"type": "string"}, "description": "影響範圍，如['個人','企業','社會','經濟']，不填則自動決定"},
            "lang": {"type": "string", "enum": ["zh-tw", "en"], "description": "語言，預設 zh-tw"}
        }, "required": ["event"]}
    },
    {
        "name": "scenario_planning",
        "description": "情境規劃：針對一個問題或決策，設計樂觀、基準、悲觀三種未來情境，並分析各情境的發生條件與因應方式。當用戶要預測未來、要規劃應對方案時使用。",
        "input_schema": {"type": "object", "properties": {
            "topic": {"type": "string", "description": "要規劃情境的主題或問題"},
            "horizon": {"type": "string", "description": "時間範圍，如1年、5年、10年（選填）"},
            "lang": {"type": "string", "enum": ["zh-tw", "en"], "description": "語言，預設 zh-tw"}
        }, "required": ["topic"]}
    },
    {
        "name": "decision_helper",
        "description": "決策輔助：針對一個決策問題蒐集相關資訊，分析各選項的優劣，最後給出結構化建議與行動步驟。當用戶面臨選擇、問該怎麼辦、要做重要決定時使用。",
        "input_schema": {"type": "object", "properties": {
            "question": {"type": "string", "description": "決策問題，如「要不要換工作」「要不要買這支股票」"},
            "options": {"type": "array", "items": {"type": "string"}, "description": "已有的選項（選填，不填則自動搜尋）"},
            "criteria": {"type": "array", "items": {"type": "string"}, "description": "決策考量標準，如['薪資','成長空間','工作環境']（選填）"}
        }, "required": ["question"]}
    },
    {
        "name": "devil_advocate",
        "description": "魔鬼代言人：刻意從相反角度挑戰一個觀點或計劃，找出最強的反駁論點，讓思考更嚴謹。當用戶要測試自己想法、要找漏洞、要做最壞打算時使用。",
        "input_schema": {"type": "object", "properties": {
            "position": {"type": "string", "description": "要被挑戰的觀點、計劃或決策"},
            "lang": {"type": "string", "enum": ["zh-tw", "en"], "description": "語言，預設 zh-tw"}
        }, "required": ["position"]}
    },
    {
        "name": "summary_writer",
        "description": "多來源摘要：搜尋多篇文章後壓縮成精華重點，讓大量資料快速被消化。當用戶要快速了解某主題、要把長資料變短、要整理重點時使用。",
        "input_schema": {"type": "object", "properties": {
            "topic": {"type": "string", "description": "要摘要的主題或關鍵字"},
            "max_points": {"type": "integer", "description": "最多幾個重點，預設7"},
            "lang": {"type": "string", "enum": ["zh-tw", "en"], "description": "語言，預設 zh-tw"}
        }, "required": ["topic"]}
    },
    {
        "name": "key_insights",
        "description": "洞察萃取：從大量搜尋資料中篩選出最有價值的 3–5 個核心洞察，去蕪存菁。當用戶要找關鍵發現、要知道最重要的幾點、要提煉精華時使用。",
        "input_schema": {"type": "object", "properties": {
            "topic": {"type": "string", "description": "要萃取洞察的主題"},
            "count": {"type": "integer", "description": "萃取幾個洞察，預設5"},
            "lang": {"type": "string", "enum": ["zh-tw", "en"], "description": "語言，預設 zh-tw"}
        }, "required": ["topic"]}
    },
    {
        "name": "bias_detector",
        "description": "偏見偵測：分析某議題相關資訊來源的立場與意識形態偏向，標記哪些資料可能有偏見，確保觀點客觀。當用戶要分辨資訊是否中立、要找出媒體立場、要確認資料可信度時使用。",
        "input_schema": {"type": "object", "properties": {
            "topic": {"type": "string", "description": "要檢測偏見的議題"},
            "lang": {"type": "string", "enum": ["zh-tw", "en"], "description": "語言，預設 zh-tw"}
        }, "required": ["topic"]}
    },
    {
        "name": "second_opinion",
        "description": "多專家視角：模擬不同領域專家（經濟學家、心理學家、工程師、醫師等）對同一問題各自給出看法，呈現多元觀點。當用戶要聽不同角度意見、要知道各領域專家怎麼看時使用。",
        "input_schema": {"type": "object", "properties": {
            "question": {"type": "string", "description": "要徵詢意見的問題或議題"},
            "experts": {"type": "array", "items": {"type": "string"}, "description": "指定專家角色，如['經濟學家','心理學家','工程師']，不填則自動決定"},
            "lang": {"type": "string", "enum": ["zh-tw", "en"], "description": "語言，預設 zh-tw"}
        }, "required": ["question"]}
    },
    {
        "name": "brainstorm",
        "description": "腦力激盪：針對一個問題搜尋現有解法後，結合創意產生多元方案與新視角。當用戶要想辦法、要找解法、要創意發想時使用。",
        "input_schema": {"type": "object", "properties": {
            "problem": {"type": "string", "description": "要腦力激盪的問題或挑戰"},
            "count": {"type": "integer", "description": "產生幾個方案，預設8"},
            "style": {"type": "string", "enum": ["實用", "創意", "顛覆"], "description": "發想風格，預設實用"},
            "lang": {"type": "string", "enum": ["zh-tw", "en"], "description": "語言，預設 zh-tw"}
        }, "required": ["problem"]}
    },
    {
        "name": "benchmark_analysis",
        "description": "標竿分析：找出該領域最佳實踐案例，與目標對象比較後給出具體改進方向。當用戶要學習最佳做法、要知道業界標準、要找對標對象時使用。",
        "input_schema": {"type": "object", "properties": {
            "subject": {"type": "string", "description": "要進行標竿分析的對象（公司、產品、做法等）"},
            "industry": {"type": "string", "description": "所屬產業或領域（選填）"},
            "lang": {"type": "string", "enum": ["zh-tw", "en"], "description": "語言，預設 zh-tw"}
        }, "required": ["subject"]}
    },
    {
        "name": "steel_man",
        "description": "鋼人論證：先把對方觀點強化到最強版本後，再表達自己的立場與反駁。比魔鬼代言人更進階，避免打稻草人，讓論述更有說服力。當用戶要深度回應某觀點、要理性辯論、要展現思想深度時使用。",
        "input_schema": {"type": "object", "properties": {
            "opposing_view": {"type": "string", "description": "要被鋼人化的對立觀點"},
            "own_position": {"type": "string", "description": "自己的立場（選填，不填則給出中立分析）"},
            "lang": {"type": "string", "enum": ["zh-tw", "en"], "description": "語言，預設 zh-tw"}
        }, "required": ["opposing_view"]}
    },
    {
        "name": "socratic_questioning",
        "description": "蘇格拉底式提問：針對一個主題產生一系列層層遞進的深層問題，引導更清晰的思考，而不只是給答案。當用戶要深入探索某議題、要訓練批判思考、要找到問題的核心時使用。",
        "input_schema": {"type": "object", "properties": {
            "topic": {"type": "string", "description": "要深入探索的主題或議題"},
            "depth": {"type": "integer", "description": "提問層數，預設5"},
            "lang": {"type": "string", "enum": ["zh-tw", "en"], "description": "語言，預設 zh-tw"}
        }, "required": ["topic"]}
    },
    {
        "name": "analogy_maker",
        "description": "類比說明：把複雜概念、技術原理或抽象看法，用生活化的類比讓人一聽就懂。當用戶要解釋難懂的東西、要讓想法更有說服力、要找到好的比喻時使用。",
        "input_schema": {"type": "object", "properties": {
            "concept": {"type": "string", "description": "要類比說明的複雜概念或想法"},
            "audience": {"type": "string", "description": "目標受眾，如一般大眾、小學生、商業人士（選填）"},
            "count": {"type": "integer", "description": "產生幾個類比，預設3"},
            "lang": {"type": "string", "enum": ["zh-tw", "en"], "description": "語言，預設 zh-tw"}
        }, "required": ["concept"]}
    },
    {
        "name": "narrative_builder",
        "description": "敘事架構：把資料與觀點包裝成「問題→衝突→洞察→結論」的有力故事結構，讓想法更有說服力與感染力。當用戶要寫有說服力的文章、要讓報告更吸引人時使用。",
        "input_schema": {"type": "object", "properties": {
            "topic": {"type": "string", "description": "要建構敘事的主題"},
            "key_message": {"type": "string", "description": "核心訊息或結論（選填）"},
            "audience": {"type": "string", "description": "目標受眾（選填）"},
            "lang": {"type": "string", "enum": ["zh-tw", "en"], "description": "語言，預設 zh-tw"}
        }, "required": ["topic"]}
    },
    {
        "name": "critique_writer",
        "description": "批判性評析：對一篇文章、政策、作品、計劃或想法進行有深度的評論，指出優點、盲點、假設前提與改進方向。當用戶要評論某事物、要寫書評影評、要深度分析一個計劃時使用。",
        "input_schema": {"type": "object", "properties": {
            "subject": {"type": "string", "description": "要評析的對象（作品名、政策名、想法描述）"},
            "type": {"type": "string", "enum": ["文章", "政策", "作品", "計劃", "觀點", "其他"], "description": "評析類型，預設觀點"},
            "lang": {"type": "string", "enum": ["zh-tw", "en"], "description": "語言，預設 zh-tw"}
        }, "required": ["subject"]}
    },
    {
        "name": "position_statement",
        "description": "立場聲明：針對爭議議題清楚表明立場，並給出「論點→證據→反駁→結論」的系統性多層次論證。當用戶要寫有說服力的意見、要公開表態、要寫辯論稿時使用。",
        "input_schema": {"type": "object", "properties": {
            "issue": {"type": "string", "description": "爭議議題"},
            "stance": {"type": "string", "enum": ["支持", "反對", "有條件支持"], "description": "立場"},
            "lang": {"type": "string", "enum": ["zh-tw", "en"], "description": "語言，預設 zh-tw"}
        }, "required": ["issue", "stance"]}
    },
    {
        "name": "ocr_click",
        "description": "OCR找字點擊：用OCR掃描指定螢幕上的所有文字，找到目標文字後自動點擊該位置。不需要模板圖，純文字定位。當用戶說「點開某某某的訊息」「點那個按鈕」「找到某個選項點一下」時使用。",
        "input_schema": {"type": "object", "properties": {
            "target_text": {"type": "string", "description": "要尋找並點擊的文字"},
            "monitor": {"type": "integer", "description": "目標螢幕編號（1/2/3），預設1"},
            "click_type": {"type": "string", "enum": ["click", "double_click", "right_click"], "description": "點擊方式，預設click"},
            "region": {"type": "array", "items": {"type": "integer"}, "description": "限定搜尋區域 [x, y, w, h]，不填則全螢幕"}
        }, "required": ["target_text"]}
    },
    {
        "name": "vision_locate",
        "description": "視覺定位：截圖後讓 Claude 視覺辨識畫面，用自然語言描述目標元素（如「Telegram左側第三個對話」「右上角關閉按鈕」），自動計算出座標並點擊。當無法用文字找到目標、需要視覺理解畫面時使用。",
        "input_schema": {"type": "object", "properties": {
            "description": {"type": "string", "description": "要找的目標元素的自然語言描述"},
            "monitor": {"type": "integer", "description": "目標螢幕編號（1/2/3），預設1"},
            "action": {"type": "string", "enum": ["click", "double_click", "right_click", "locate_only"], "description": "找到後執行的動作，預設click"},
            "region": {"type": "array", "items": {"type": "integer"}, "description": "限定截圖區域 [x, y, w, h]，不填則全螢幕"}
        }, "required": ["description"]}
    },
    {
        "name": "screen_workflow",
        "description": "螢幕工作流：把「截圖→辨識→點擊→等待→輸入→送出」串成多步驟自動化序列，一個指令完成整件事。當用戶說「幫我打開Telegram回覆某某某說○○」這類需要多個步驟的操作時使用。",
        "input_schema": {"type": "object", "properties": {
            "steps": {"type": "array", "items": {
                "type": "object", "properties": {
                    "action": {"type": "string", "description": "動作類型：screenshot/ocr_click/vision_click/type/press/wait/open_app"},
                    "target": {"type": "string", "description": "目標文字或描述"},
                    "value": {"type": "string", "description": "輸入值或等待秒數"},
                    "monitor": {"type": "integer", "description": "螢幕編號"}
                }
            }, "description": "步驟清單，依序執行"}
        }, "required": ["steps"]}
    },
    {
        "name": "app_navigator",
        "description": "【優先使用】App多步驟自動操作：當用戶說以下任何一種時必須使用此工具：「從Telegram/LINE找XXX」「在Telegram跟XXX說」「幫我打開Telegram找好友XXX」「Telegram傳訊息給XXX說OOO」「螢幕X的Telegram找XXX並回覆」「LINE傳給XXX說OOO」。凡是涉及在App內找聯絡人、打開對話、發訊息、回覆訊息的複合操作，一律用這個工具，不要用 desktop_control 截圖。",
        "input_schema": {"type": "object", "properties": {
            "app": {"type": "string", "description": "目標App名稱，如Telegram、LINE、Chrome、記事本"},
            "task": {"type": "string", "description": "要執行的任務，用自然語言描述，如「打開某某某的對話」「點新增按鈕」"},
            "contact_name": {"type": "string", "description": "【Telegram/LINE必填，不可省略】要搜尋的聯絡人名稱，只填純名字，例如「巴斯」「奈絲菟米啾」，絕對不能含「找」「跟」「並打開」等動詞"},
            "input_text": {"type": "string", "description": "需要輸入的文字（如回覆內容），選填"},
            "monitor": {"type": "integer", "description": "App所在螢幕編號（1/2/3），預設1"}
        }, "required": ["app", "task"]}
    },
    {
        "name": "wait_and_click",
        "description": "等待出現後點擊：持續監控螢幕，等到指定文字或圖示出現後才執行點擊，解決非同步等待問題（如等訊息載入、等彈窗出現）。當用戶的操作需要等待某個元素出現時使用。",
        "input_schema": {"type": "object", "properties": {
            "target_text": {"type": "string", "description": "等待出現的目標文字"},
            "timeout": {"type": "integer", "description": "最長等待秒數，預設15"},
            "monitor": {"type": "integer", "description": "監控的螢幕編號（1/2/3），預設1"},
            "action_after": {"type": "string", "enum": ["click", "double_click", "none"], "description": "出現後執行的動作，預設click"}
        }, "required": ["target_text"]}
    },
    {
        "name": "drag_drop",
        "description": "拖曳操作：從一個座標拖曳到另一個座標，支援文字描述定位拖曳起點和終點。用於排序、移動視窗、拖曳上傳檔案、調整大小等操作。",
        "input_schema": {"type": "object", "properties": {
            "from_x": {"type": "integer", "description": "起點X座標（與from_text二選一）"},
            "from_y": {"type": "integer", "description": "起點Y座標"},
            "to_x": {"type": "integer", "description": "終點X座標（與to_text二選一）"},
            "to_y": {"type": "integer", "description": "終點Y座標"},
            "from_text": {"type": "string", "description": "起點目標文字（OCR定位，與from_x/y二選一）"},
            "to_text": {"type": "string", "description": "終點目標文字（OCR定位，與to_x/y二選一）"},
            "monitor": {"type": "integer", "description": "螢幕編號，預設1"},
            "duration": {"type": "number", "description": "拖曳時間秒數，預設0.5"}
        }, "required": []}
    },
    {
        "name": "read_screen",
        "description": "讀取螢幕內容：截圖後用 Vision AI 分析螢幕上有什麼，回傳詳細描述。用於「幫我看一下螢幕X現在顯示什麼」「確認有沒有出現XXX」「告訴我現在螢幕上的內容」。",
        "input_schema": {"type": "object", "properties": {
            "question": {"type": "string", "description": "要詢問螢幕的問題，如「現在顯示什麼」「有沒有出現確認對話框」，預設描述全部內容"},
            "monitor": {"type": "integer", "description": "螢幕編號（1/2/3），預設1"}
        }, "required": []}
    },
    {
        "name": "scroll_at",
        "description": "在螢幕指定位置滾動：支援所有螢幕包含螢幕2。可用座標或文字描述定位滾動位置。用於「螢幕X往下滾」「滾動到看到XXX」「在Telegram聊天列表往上滾」。",
        "input_schema": {"type": "object", "properties": {
            "direction": {"type": "string", "enum": ["up", "down"], "description": "滾動方向，預設down"},
            "amount": {"type": "integer", "description": "滾動格數，預設3"},
            "x": {"type": "integer", "description": "滾動位置X（相對螢幕左上角），與description二選一"},
            "y": {"type": "integer", "description": "滾動位置Y（相對螢幕左上角）"},
            "monitor": {"type": "integer", "description": "螢幕編號（1/2/3），預設1"},
            "description": {"type": "string", "description": "用文字描述滾動目標區域，如「Telegram左側聊天列表」"}
        }, "required": []}
    },
    {
        "name": "window_manager",
        "description": "視窗管理：列出所有開啟的視窗、切換視窗到前景、最大化/最小化/關閉視窗。用於「列出所有視窗」「把Telegram切到前景」「最大化記事本」「關閉Chrome」。",
        "input_schema": {"type": "object", "properties": {
            "action": {"type": "string", "enum": ["list", "focus", "maximize", "minimize", "close"],
                       "description": "動作：list列出所有視窗 / focus切換到前景 / maximize最大化 / minimize最小化 / close關閉"},
            "window_name": {"type": "string", "description": "視窗名稱關鍵字（list動作不需要填）"}
        }, "required": ["action"]}
    },
    {
        "name": "get_candlestick_chart",
        "description": "產生股票 K 線圖並傳送圖片，包含成交量、MA均線、自動辨識型態（頭肩頂、雙底、突破等）。當用戶要看 K 線圖、走勢圖時使用。",
        "input_schema": {
            "type": "object",
            "properties": {
                "symbol": {"type": "string", "description": "股票代號"},
                "period": {"type": "string", "enum": ["1mo", "3mo", "6mo", "1y"], "description": "期間，預設 3mo"}
            },
            "required": ["symbol"]
        }
    },
    {
        "name": "compare_stocks",
        "description": "比較多支股票的表現：漲跌幅、本益比、ROE、市值等關鍵指標並排比較。當用戶問哪支比較好、要比較兩支股票時使用。",
        "input_schema": {
            "type": "object",
            "properties": {
                "symbols": {"type": "array", "items": {"type": "string"}, "description": "要比較的股票代號列表，最多5支"},
                "metrics": {"type": "array", "items": {"type": "string", "enum": ["price", "pe", "roe", "margin", "growth", "all"]}, "description": "比較項目，預設 all"}
            },
            "required": ["symbols"]
        }
    },
    {
        "name": "ptt_search",
        "description": "搜尋 PTT 熱門文章與留言，了解台灣民間對某話題的真實討論與輿論。當用戶問台灣人怎麼看、PTT 怎麼說、台灣社會觀感時使用。",
        "input_schema": {
            "type": "object",
            "properties": {
                "board": {"type": "string", "description": "看板名稱，如 Gossiping（八卦）、Stock（股票）、HatePolitics（政黑）、Baseball（棒球）、Tech_Job，預設 Gossiping"},
                "keyword": {"type": "string", "description": "搜尋關鍵字"},
                "count": {"type": "integer", "description": "幾篇文章，預設 5"}
            },
            "required": ["keyword"]
        }
    },
    {
        "name": "multi_perspective",
        "description": "針對同一話題，從正方、反方、中立三個角度搜尋觀點並整合分析。當用戶問某事好不好、某人評價、某政策對不對、要深度分析一個爭議話題時使用。",
        "input_schema": {
            "type": "object",
            "properties": {
                "topic": {"type": "string", "description": "要分析的話題、人物或事件"},
                "lang": {"type": "string", "description": "搜尋語言，預設 zh-tw", "enum": ["zh-tw", "en"]}
            },
            "required": ["topic"]
        }
    },
    {
        "name": "google_trends",
        "description": "查詢 Google Trends 熱度，了解某話題在一段時間內的搜尋趨勢與熱門程度。當用戶問某話題熱不熱、是否在發酵時使用。",
        "input_schema": {
            "type": "object",
            "properties": {
                "keywords": {"type": "array", "items": {"type": "string"}, "description": "要比較的關鍵字，最多5個"},
                "timeframe": {"type": "string", "description": "時間範圍，預設 today 3-m（近3月）", "enum": ["today 1-m", "today 3-m", "today 12-m", "today 5-y"]},
                "geo": {"type": "string", "description": "地區，預設 TW（台灣），可用 US、CN、HK、JP"}
            },
            "required": ["keywords"]
        }
    },
    {
        "name": "read_webpage",
        "description": "讀取網頁完整內容並分析。當用戶貼網址、要分析某篇文章、要深度了解某個網頁內容時使用。比搜尋更深入，直接讀取全文。",
        "input_schema": {
            "type": "object",
            "properties": {
                "url": {"type": "string", "description": "要讀取的網址"},
                "max_chars": {"type": "integer", "description": "最多擷取幾個字元，預設 3000"}
            },
            "required": ["url"]
        }
    },
    {
        "name": "wikipedia_search",
        "description": "查詢 Wikipedia 百科，取得人物、事件、組織、地點的背景知識。當用戶問某人是誰、某事件的來龍去脈、某組織的背景時使用。",
        "input_schema": {
            "type": "object",
            "properties": {
                "query": {"type": "string", "description": "要查詢的人物、事件或主題"},
                "lang": {"type": "string", "description": "語言版本，預設 zh（中文），可用 en（英文）", "enum": ["zh", "en", "ja"]}
            },
            "required": ["query"]
        }
    },
    {
        "name": "news_search",
        "description": "搜尋任何主題的最新新聞（非財經專用）。當用戶問某人最近的新聞、某事件最新動態、某話題的報導時使用。",
        "input_schema": {
            "type": "object",
            "properties": {
                "query": {"type": "string", "description": "搜尋關鍵字"},
                "lang": {"type": "string", "description": "語言，預設 zh-TW", "enum": ["zh-TW", "zh-CN", "en-US"]},
                "count": {"type": "integer", "description": "幾則新聞，預設 6"}
            },
            "required": ["query"]
        }
    },
    {
        "name": "youtube_summary",
        "description": "擷取 YouTube 影片字幕並分析內容、觀點、重點摘要。當用戶貼 YouTube 連結要了解影片說什麼時使用。",
        "input_schema": {
            "type": "object",
            "properties": {
                "url": {"type": "string", "description": "YouTube 影片網址或影片 ID"}
            },
            "required": ["url"]
        }
    },
    {
        "name": "analyze_pdf",
        "description": "讀取並分析 PDF 文件內容。當用戶傳來 PDF 檔案路徑或要分析報告、研究文件時使用。",
        "input_schema": {
            "type": "object",
            "properties": {
                "path": {"type": "string", "description": "PDF 檔案的完整路徑"},
                "max_chars": {"type": "integer", "description": "最多擷取幾個字元，預設 4000"}
            },
            "required": ["path"]
        }
    },
    {
        "name": "ddg_search",
        "description": "DuckDuckGo 快速網路搜尋，適合搜尋中文內容、中國相關話題、最新時事、任何需要查資料的問題。比 osint_search 更快更穩。當用戶問不知道答案的問題、需要查最新資訊、詢問中國相關話題時優先使用此工具。",
        "input_schema": {
            "type": "object",
            "properties": {
                "query": {"type": "string", "description": "搜尋關鍵字，可用中文或英文"},
                "region": {"type": "string", "description": "搜尋地區語言，預設 zh-tw（繁中），可用 zh-cn（簡中）、us-en（英文）", "enum": ["zh-tw", "zh-cn", "us-en", "wt-wt"]},
                "max_results": {"type": "integer", "description": "最多幾筆結果，預設 5，最多 10"}
            },
            "required": ["query"]
        }
    },
    {
        "name": "get_stock_advanced",
        "description": "股票進階技術分析：MACD、布林通道（BB）、KD隨機指標。當用戶問 MACD、布林通道、KD、進階技術分析時使用。",
        "input_schema": {
            "type": "object",
            "properties": {
                "symbol": {"type": "string", "description": "股票代號，同 get_stock"},
                "indicators": {"type": "array", "items": {"type": "string", "enum": ["macd", "bb", "kd"]}, "description": "要計算的指標，可多選，預設全部"}
            },
            "required": ["symbol"]
        }
    },
    {
        "name": "get_macro",
        "description": "查詢總體經濟指標：美國CPI通膨、失業率、GDP成長率、短期利率。當用戶問總經、通膨、Fed利率、經濟數據時使用。",
        "input_schema": {
            "type": "object",
            "properties": {
                "indicator": {"type": "string", "enum": ["cpi", "unemployment", "fed_rate", "gdp", "nonfarm"], "description": "cpi=通膨年增率, unemployment=失業率, fed_rate=短期利率, gdp=GDP成長率, nonfarm=非農就業說明"}
            },
            "required": ["indicator"]
        }
    },
    {
        "name": "portfolio",
        "description": "個人投資組合管理：新增持股、刪除持股、查看總損益與報酬率。當用戶要管理自己的股票持倉、追蹤損益時使用。",
        "input_schema": {
            "type": "object",
            "properties": {
                "action": {"type": "string", "enum": ["add", "remove", "view", "clear"], "description": "add=新增持股, remove=刪除, view=查看損益, clear=清空"},
                "chat_id": {"type": "integer", "description": "聊天室 ID"},
                "symbol": {"type": "string", "description": "股票代號（add/remove 用）"},
                "shares": {"type": "number", "description": "持有數量（add 用）"},
                "cost": {"type": "number", "description": "買入均價（add 用）"}
            },
            "required": ["action"]
        }
    },
    {
        "name": "desktop_control",
        "description": "控制電腦桌面基本操作，支援多螢幕。【強制】當用戶說「打開XXX」「開啟XXX」「執行XXX」時，你必須呼叫此工具的 open_app action，不能只用文字回覆。支援：(1)截圖 (2)點擊座標 (3)輸入文字 (4)按鍵 (5)開啟程式 (6)查看螢幕。如果任務是「在某App找某人並操作」請改用 app_navigator。",
        "input_schema": {
            "type": "object",
            "properties": {
                "action": {
                    "type": "string",
                    "enum": ["screenshot", "click", "double_click", "right_click", "move", "type", "press_key", "open_app", "scroll", "list_monitors"],
                    "description": "要執行的動作。list_monitors=列出所有螢幕資訊"
                },
                "x": {"type": "integer", "description": "滑鼠 X 座標，若指定 monitor 則為該螢幕內的相對座標"},
                "y": {"type": "integer", "description": "滑鼠 Y 座標，若指定 monitor 則為該螢幕內的相對座標"},
                "monitor": {"type": "integer", "description": "目標螢幕編號（1=左、2=右、3=中，不填則為全域座標）。screenshot 時只截該螢幕"},
                "text": {"type": "string", "description": "要輸入的文字或按鍵名稱（type/press_key 時使用）"},
                "app": {"type": "string", "description": "要開啟的程式名稱或路徑（open_app 時使用）"},
                "direction": {"type": "string", "enum": ["up", "down"], "description": "滾動方向（scroll 時使用）"},
                "amount": {"type": "integer", "description": "滾動格數（scroll 時使用，預設 3）"}
            },
            "required": ["action"]
        }
    },
    {
        "name": "screen_vision",
        "description": "截圖並分析畫面內容。當用戶問「電腦現在在做什麼」、「幫我看一下螢幕」、「截圖分析」時使用。",
        "input_schema": {
            "type": "object",
            "properties": {
                "question": {"type": "string", "description": "要問關於畫面的問題，預設為「請描述畫面上有什麼」"}
            },
            "required": []
        }
    },
    {
        "name": "tg_auto_reply",
        "description": "開啟/停止 Telegram 自動回覆監控。監控螢幕2的 Telegram 對話，偵測到對方新訊息時自動用小牛馬風格回覆。用戶說「開啟自動回覆」「監控聊天」「幫我回訊息」時使用。當用戶說「到幾點」「時間到XX:XX」時，必須用 stop_time 參數帶入結束時間。",
        "input_schema": {
            "type": "object",
            "properties": {
                "duration_minutes": {"type": "number", "description": "監控持續時間（分鐘），預設 30。如果用戶指定了結束時間則不需要此參數"},
                "stop_time": {"type": "string", "description": "監控結束時間，格式 HH:MM（如 '15:15'）。用戶說「到15:15」「時間到15:15」時必須填入此參數，優先於 duration_minutes"},
                "action": {"type": "string", "description": "start 開啟 / stop 停止，預設 start"}
            },
            "required": []
        }
    },
    {
        "name": "find_image_on_screen",
        "description": "在螢幕上找到指定圖片檔案的位置，回傳座標。用於視覺定位按鈕或元素。",
        "input_schema": {
            "type": "object",
            "properties": {
                "template_path": {"type": "string", "description": "要尋找的圖片檔案路徑"},
                "confidence": {"type": "number", "description": "匹配信心度 0~1，預設 0.8"}
            },
            "required": ["template_path"]
        }
    },
    {
        "name": "browser_control",
        "description": "控制瀏覽器自動化，用於互動操作：點擊、填表、截圖。若只是要讀取網頁文字內容，請用 web_scrape 而非此工具。",
        "input_schema": {
            "type": "object",
            "properties": {
                "action": {
                    "type": "string",
                    "enum": ["open", "click", "type", "get_text", "screenshot", "goto", "close"],
                    "description": "open=開啟網址, click=點擊選擇器, type=輸入文字, get_text=取得文字, screenshot=截圖, goto=前往網址, close=關閉"
                },
                "url": {"type": "string", "description": "網址（open/goto 時使用）"},
                "selector": {"type": "string", "description": "CSS 選擇器（click/type/get_text 時使用）"},
                "text": {"type": "string", "description": "要輸入的文字（type 時使用）"}
            },
            "required": ["action"]
        }
    },
    {
        "name": "window_control",
        "description": "管理視窗。列出、切換、最大化、最小化、關閉視窗。",
        "input_schema": {
            "type": "object",
            "properties": {
                "action": {"type": "string", "enum": ["list","focus","close","minimize","maximize"]},
                "keyword": {"type": "string", "description": "視窗標題關鍵字"}
            },
            "required": ["action"]
        }
    },
    {
        "name": "hotkey",
        "description": "執行鍵盤組合鍵，如 ctrl+c、alt+tab、win+d。",
        "input_schema": {
            "type": "object",
            "properties": {
                "keys": {"type": "string", "description": "組合鍵，用加號連接，例如 ctrl+c 或 alt+tab"}
            },
            "required": ["keys"]
        }
    },
    {
        "name": "clipboard",
        "description": "讀取或寫入剪貼簿內容。",
        "input_schema": {
            "type": "object",
            "properties": {
                "action": {"type": "string", "enum": ["get","set"]},
                "text": {"type": "string", "description": "要寫入的文字（set 時使用）"}
            },
            "required": ["action"]
        }
    },
    {
        "name": "file_system",
        "description": "操作檔案與資料夾：列出、讀取、寫入、刪除、複製、移動、搜尋。",
        "input_schema": {
            "type": "object",
            "properties": {
                "action": {"type": "string", "enum": ["list","read","write","delete","copy","move","search"]},
                "path": {"type": "string", "description": "檔案或資料夾路徑"},
                "dest": {"type": "string", "description": "目標路徑（copy/move 時使用）"},
                "content": {"type": "string", "description": "寫入內容（write 時使用）"},
                "keyword": {"type": "string", "description": "搜尋關鍵字（search 時使用）"}
            },
            "required": ["action"]
        }
    },
    {
        "name": "system_monitor",
        "description": "查看系統資源使用狀況（CPU、記憶體、磁碟）或管理程序。",
        "input_schema": {
            "type": "object",
            "properties": {
                "action": {"type": "string", "enum": ["info","process_list","kill"]},
                "target": {"type": "string", "description": "程序名稱或 PID（kill 時使用）"}
            },
            "required": ["action"]
        }
    },
    {
        "name": "notify",
        "description": "發送 Windows 桌面通知彈出視窗。",
        "input_schema": {
            "type": "object",
            "properties": {
                "title": {"type": "string"},
                "message": {"type": "string"}
            },
            "required": ["title","message"]
        }
    },
    {
        "name": "tts",
        "description": "文字轉語音，讓電腦朗讀指定文字。",
        "input_schema": {
            "type": "object",
            "properties": {
                "text": {"type": "string"}
            },
            "required": ["text"]
        }
    },
    {
        "name": "ai_plan",
        "description": "AI 自動規劃並執行多步驟任務。用戶給一個目標，自動拆解成步驟並執行。",
        "input_schema": {
            "type": "object",
            "properties": {
                "goal": {"type": "string", "description": "要達成的目標"}
            },
            "required": ["goal"]
        }
    },
    {
        "name": "screen_stream",
        "description": "螢幕串流，持續截圖並傳送給用戶。",
        "input_schema": {
            "type": "object",
            "properties": {
                "duration": {"type": "integer", "description": "持續秒數，預設 10"},
                "interval": {"type": "number", "description": "截圖間隔秒數，預設 2"}
            },
            "required": []
        }
    },
    {
        "name": "drag",
        "description": "拖曳滑鼠從一個位置到另一個位置。",
        "input_schema": {
            "type": "object",
            "properties": {
                "x1": {"type": "integer"}, "y1": {"type": "integer"},
                "x2": {"type": "integer"}, "y2": {"type": "integer"},
                "duration": {"type": "number", "description": "拖曳時間，預設 0.5 秒"}
            },
            "required": ["x1","y1","x2","y2"]
        }
    },
    {
        "name": "power_control",
        "description": "電源管理：睡眠、重開機、關機。",
        "input_schema": {
            "type": "object",
            "properties": {
                "action": {"type": "string", "enum": ["sleep","restart","shutdown"]}
            },
            "required": ["action"]
        }
    },
    {
        "name": "virtual_desktop",
        "description": "切換或建立 Windows 虛擬桌面。",
        "input_schema": {
            "type": "object",
            "properties": {
                "action": {"type": "string", "enum": ["left","right","new"]}
            },
            "required": ["action"]
        }
    },
    {
        "name": "bluetooth",
        "description": "掃描或連線藍牙裝置。",
        "input_schema": {
            "type": "object",
            "properties": {
                "action": {"type": "string", "enum": ["scan","connect"]},
                "mac": {"type": "string", "description": "藍牙 MAC 位址（connect 時使用）"}
            },
            "required": ["action"]
        }
    },
    {
        "name": "stt",
        "description": "語音辨識，錄製麥克風聲音並轉為文字。",
        "input_schema": {"type": "object", "properties": {}, "required": []}
    },
    {
        "name": "ocr",
        "description": "OCR 文字辨識，截取畫面或指定圖片，辨識其中的文字內容。",
        "input_schema": {
            "type": "object",
            "properties": {
                "image_path": {"type": "string", "description": "圖片路徑（選填，不填則截圖）"}
            },
            "required": []
        }
    },
    {
        "name": "workflow",
        "description": "執行或儲存自動化工作流程（多步驟串接）。",
        "input_schema": {
            "type": "object",
            "properties": {
                "action": {"type": "string", "enum": ["run", "save", "list"]},
                "name": {"type": "string", "description": "流程名稱"},
                "steps": {"type": "string", "description": "JSON 格式的步驟（save 時使用）"}
            },
            "required": ["action"]
        }
    },
    {
        "name": "screen_watch",
        "description": "監控螢幕，當出現指定圖片時自動執行指令。",
        "input_schema": {
            "type": "object",
            "properties": {
                "template_path": {"type": "string"},
                "command": {"type": "string"},
                "timeout": {"type": "number", "description": "等待秒數，預設 60"}
            },
            "required": ["template_path", "command"]
        }
    },
    {
        "name": "file_transfer",
        "description": "壓縮/解壓縮檔案或下載網路檔案。",
        "input_schema": {
            "type": "object",
            "properties": {
                "action": {"type": "string", "enum": ["zip", "unzip", "download"]},
                "source": {"type": "string", "description": "來源路徑或網址"},
                "dest": {"type": "string", "description": "目標路徑"}
            },
            "required": ["action", "source"]
        }
    },
    {
        "name": "send_email",
        "description": "發送電子郵件。需要 .env 設定 SMTP_USER 和 SMTP_PASS。",
        "input_schema": {
            "type": "object",
            "properties": {
                "to": {"type": "string"},
                "subject": {"type": "string"},
                "body": {"type": "string"}
            },
            "required": ["to","subject","body"]
        }
    },
    {
        "name": "long_term_memory",
        "description": "管理長期記憶。當用戶說了重要的事（偏好、習慣、重要資訊、個人資料）或要求記住某件事時，主動儲存。當用戶問你記不記得某件事時，先查詢記憶再回答。",
        "input_schema": {
            "type": "object",
            "properties": {
                "action": {
                    "type": "string",
                    "enum": ["save", "list", "delete"],
                    "description": "save=儲存新記憶, list=列出所有記憶, delete=刪除指定記憶"
                },
                "content": {"type": "string", "description": "要儲存的記憶內容（save 時使用）"},
                "memory_id": {"type": "integer", "description": "要刪除的記憶 ID（delete 時使用）"}
            },
            "required": ["action"]
        }
    },
    {
        "name": "manage_schedule",
        "description": "管理 Windows 排程任務與 Bot 狀態。當用戶要查看排程、新增排程、刪除排程、重啟 bot、查看 bot 是否在跑時使用此工具。",
        "input_schema": {
            "type": "object",
            "properties": {
                "action": {
                    "type": "string",
                    "enum": ["list", "add", "delete", "bot_status", "bot_restart"],
                    "description": "list=列出所有排程, add=新增排程, delete=刪除排程, bot_status=查看bot狀態, bot_restart=重啟bot"
                },
                "name": {"type": "string", "description": "排程任務名稱（add/delete 時使用）"},
                "time": {"type": "string", "description": "執行時間，格式 HH:MM（add 時使用）"},
                "script": {"type": "string", "description": "腳本完整路徑（add 時使用）"}
            },
            "required": ["action"]
        }
    },
    {
        "name": "volume",
        "description": "控制系統音量：查詢、設定音量大小、靜音/取消靜音。",
        "input_schema": {
            "type": "object",
            "properties": {
                "action": {"type": "string", "enum": ["get","set","mute","unmute"]},
                "level": {"type": "integer", "description": "音量 0-100（set 使用）"}
            },
            "required": ["action"]
        }
    },
    {
        "name": "display",
        "description": "控制螢幕亮度或查詢解析度。",
        "input_schema": {
            "type": "object",
            "properties": {
                "action": {"type": "string", "enum": ["brightness_get","brightness_set","resolution"]},
                "level": {"type": "integer", "description": "亮度 0-100（brightness_set 使用）"}
            },
            "required": ["action"]
        }
    },
    {
        "name": "media",
        "description": "控制媒體播放：播放/暫停、上下首、停止、音量、切換音訊裝置。",
        "input_schema": {
            "type": "object",
            "properties": {
                "action": {"type": "string", "enum": ["play_pause","next","prev","stop","volume_up","volume_down","mute","list_devices","switch_device"]},
                "device_name": {"type": "string", "description": "裝置名稱（switch_device 使用）"}
            },
            "required": ["action"]
        }
    },
    {
        "name": "software",
        "description": "管理已安裝軟體：列出、安裝、卸載。",
        "input_schema": {
            "type": "object",
            "properties": {
                "action": {"type": "string", "enum": ["list","install","uninstall"]},
                "name": {"type": "string", "description": "軟體名稱或 winget ID"},
                "keyword": {"type": "string", "description": "搜尋關鍵字（list 使用）"}
            },
            "required": ["action"]
        }
    },
    {
        "name": "startup",
        "description": "管理開機自啟動程式：列出、新增、移除。",
        "input_schema": {
            "type": "object",
            "properties": {
                "action": {"type": "string", "enum": ["list","add","remove"]},
                "name": {"type": "string"},
                "command": {"type": "string", "description": "啟動指令（add 使用）"}
            },
            "required": ["action"]
        }
    },
    {
        "name": "env_var",
        "description": "查詢或設定 Windows 環境變數。",
        "input_schema": {
            "type": "object",
            "properties": {
                "action": {"type": "string", "enum": ["get","set"]},
                "name": {"type": "string"},
                "value": {"type": "string", "description": "值（set 使用）"},
                "permanent": {"type": "boolean", "description": "是否永久設定"}
            },
            "required": ["action"]
        }
    },
    {
        "name": "user_account",
        "description": "管理 Windows 使用者帳戶：列出、建立、刪除。",
        "input_schema": {
            "type": "object",
            "properties": {
                "action": {"type": "string", "enum": ["list","create","delete"]},
                "username": {"type": "string"},
                "password": {"type": "string", "description": "密碼（create 使用）"}
            },
            "required": ["action"]
        }
    },
    {
        "name": "windows_update",
        "description": "管理 Windows 系統更新：查看可用更新、執行更新。",
        "input_schema": {
            "type": "object",
            "properties": {
                "action": {"type": "string", "enum": ["list","install","check"]}
            },
            "required": ["action"]
        }
    },
    {
        "name": "device_manager",
        "description": "裝置管理員：列出硬體裝置、啟用或停用指定裝置。",
        "input_schema": {
            "type": "object",
            "properties": {
                "action": {"type": "string", "enum": ["list","enable","disable"]},
                "name": {"type": "string", "description": "裝置名稱關鍵字"},
                "keyword": {"type": "string", "description": "搜尋關鍵字（list 使用）"}
            },
            "required": ["action"]
        }
    },
    {
        "name": "network_config",
        "description": "網路進階設定：網路卡啟停、DNS設定、靜態IP、Hosts編輯、流量監控。",
        "input_schema": {
            "type": "object",
            "properties": {
                "action": {"type": "string", "enum": ["adapter_list","adapter_enable","adapter_disable","dns_get","dns_set","ip_set","hosts_list","hosts_add","hosts_remove","traffic"]},
                "name": {"type": "string", "description": "介面/裝置名稱"},
                "ip": {"type": "string"},
                "dns1": {"type": "string"},
                "dns2": {"type": "string"},
                "domain": {"type": "string"},
                "duration": {"type": "integer", "description": "監控秒數（traffic 使用）"}
            },
            "required": ["action"]
        }
    },
    {
        "name": "automation",
        "description": "進階自動化：條件式觸發執行、多視窗排列、指定區域OCR、指定視窗截圖。",
        "input_schema": {
            "type": "object",
            "properties": {
                "action": {"type": "string", "enum": ["if_then","window_arrange","region_ocr","window_screenshot"]},
                "condition_type": {"type": "string", "enum": ["cpu_above","mem_above","file_exists","time_is","process_running"], "description": "條件類型（if_then 使用）"},
                "condition_value": {"type": "string", "description": "條件值"},
                "command": {"type": "string", "description": "觸發時執行的指令"},
                "duration": {"type": "number", "description": "監控秒數"},
                "layout": {"type": "string", "enum": ["side_by_side","quad","stack","maximize_all"], "description": "視窗排列方式"},
                "x": {"type": "integer"}, "y": {"type": "integer"},
                "w": {"type": "integer"}, "h": {"type": "integer"},
                "keyword": {"type": "string", "description": "視窗標題關鍵字"},
                "output": {"type": "string"}
            },
            "required": ["action"]
        }
    },
    {
        "name": "firewall",
        "description": "管理 Windows 防火牆規則：列出、新增、刪除規則，啟用/停用防火牆。",
        "input_schema": {
            "type": "object",
            "properties": {
                "action": {"type": "string", "enum": ["list","add","remove","enable","disable","status"]},
                "name": {"type": "string", "description": "規則名稱"},
                "port": {"type": "integer", "description": "連接埠（add 使用）"},
                "protocol": {"type": "string", "enum": ["TCP","UDP","Any"], "description": "協定"},
                "direction": {"type": "string", "enum": ["Inbound","Outbound"], "description": "方向"}
            },
            "required": ["action"]
        }
    },
    {
        "name": "process_mgr",
        "description": "進階程序管理：列出/搜尋程序、強制終止、調整 CPU 優先權。",
        "input_schema": {
            "type": "object",
            "properties": {
                "action": {"type": "string", "enum": ["list","kill","priority","search"]},
                "name": {"type": "string", "description": "程序名稱關鍵字"},
                "pid": {"type": "integer", "description": "程序 ID（kill/priority 使用）"},
                "level": {"type": "string", "enum": ["realtime","high","above_normal","normal","below_normal","idle"], "description": "優先權等級"}
            },
            "required": ["action"]
        }
    },
    {
        "name": "power_plan",
        "description": "管理 Windows 電源計畫：查詢目前計畫、切換高效能/省電/均衡模式。",
        "input_schema": {
            "type": "object",
            "properties": {
                "action": {"type": "string", "enum": ["get","set","list"]},
                "plan": {"type": "string", "enum": ["balanced","high_performance","power_saver"], "description": "電源計畫（set 使用）"}
            },
            "required": ["action"]
        }
    },
    {
        "name": "event_log",
        "description": "讀取 Windows 事件記錄：系統/應用程式/安全性日誌中的錯誤、警告。",
        "input_schema": {
            "type": "object",
            "properties": {
                "log": {"type": "string", "enum": ["System","Application","Security"], "description": "日誌類型"},
                "level": {"type": "string", "enum": ["Error","Warning","Information","All"], "description": "事件等級"},
                "count": {"type": "integer", "description": "筆數，預設 10"}
            },
            "required": ["log"]
        }
    },
    {
        "name": "datetime_config",
        "description": "設定系統時間或時區、同步網路時間。",
        "input_schema": {
            "type": "object",
            "properties": {
                "action": {"type": "string", "enum": ["get","sync","set_timezone","set_time"]},
                "timezone": {"type": "string", "description": "時區名稱，例如 'Taipei Standard Time'"},
                "datetime": {"type": "string", "description": "日期時間，格式 YYYY-MM-DD HH:MM:SS（set_time 使用）"}
            },
            "required": ["action"]
        }
    },
    {
        "name": "ui_auto",
        "description": "UI 自動化：用 Windows Accessibility API 點擊按鈕、讀取視窗文字、填入表單，不需截圖。",
        "input_schema": {
            "type": "object",
            "properties": {
                "action": {"type": "string", "enum": ["find","click","type","read","get_windows"]},
                "window": {"type": "string", "description": "視窗標題關鍵字"},
                "control": {"type": "string", "description": "控制項標題/類型關鍵字"},
                "text": {"type": "string", "description": "輸入文字（type 使用）"}
            },
            "required": ["action"]
        }
    },
    {
        "name": "macro",
        "description": "鍵盤滑鼠巨集：開始錄製操作、停止錄製、回放已錄製的巨集。",
        "input_schema": {
            "type": "object",
            "properties": {
                "action": {"type": "string", "enum": ["record_start","record_stop","play","list","delete"]},
                "name": {"type": "string", "description": "巨集名稱"},
                "repeat": {"type": "integer", "description": "回放次數，預設 1"},
                "duration": {"type": "number", "description": "錄製秒數（record_start 使用）"}
            },
            "required": ["action"]
        }
    },
    {
        "name": "color_pick",
        "description": "取得螢幕指定座標的顏色值（RGB 和 HEX），或截取特定區域的主要顏色。",
        "input_schema": {
            "type": "object",
            "properties": {
                "action": {"type": "string", "enum": ["pick","dominant"]},
                "x": {"type": "integer", "description": "X 座標"},
                "y": {"type": "integer", "description": "Y 座標"},
                "region_w": {"type": "integer", "description": "區域寬度（dominant 使用）"},
                "region_h": {"type": "integer", "description": "區域高度（dominant 使用）"}
            },
            "required": ["action","x","y"]
        }
    },
    {
        "name": "webcam",
        "description": "使用 USB 攝影機拍照或錄製短片。",
        "input_schema": {
            "type": "object",
            "properties": {
                "action": {"type": "string", "enum": ["photo","video","list"]},
                "duration": {"type": "number", "description": "錄影秒數（video 使用）"},
                "output": {"type": "string", "description": "輸出檔案路徑"},
                "device": {"type": "integer", "description": "攝影機編號，預設 0"}
            },
            "required": ["action"]
        }
    },
    {
        "name": "multi_monitor",
        "description": "多螢幕管理：列出螢幕資訊、設定主螢幕、切換延伸/複製模式、移動視窗到指定螢幕。",
        "input_schema": {
            "type": "object",
            "properties": {
                "action": {"type": "string", "enum": ["list","set_primary","extend","clone","move_window"]},
                "monitor": {"type": "integer", "description": "螢幕編號（從 1 開始）"},
                "window": {"type": "string", "description": "視窗標題關鍵字（move_window 使用）"}
            },
            "required": ["action"]
        }
    },
    {
        "name": "printer",
        "description": "印表機管理：列出印表機、列印文件、查看/清除列印佇列。",
        "input_schema": {
            "type": "object",
            "properties": {
                "action": {"type": "string", "enum": ["list","print","queue","clear_queue","set_default"]},
                "path": {"type": "string", "description": "檔案路徑（print 使用）"},
                "printer_name": {"type": "string", "description": "印表機名稱"}
            },
            "required": ["action"]
        }
    },
    {
        "name": "wifi",
        "description": "Wi-Fi 管理：掃描附近網路、連線、斷線、查看目前連線、顯示已儲存密碼。",
        "input_schema": {
            "type": "object",
            "properties": {
                "action": {"type": "string", "enum": ["scan","connect","disconnect","status","saved","password"]},
                "ssid": {"type": "string", "description": "Wi-Fi 名稱（connect 使用）"},
                "password": {"type": "string", "description": "Wi-Fi 密碼（connect 使用）"}
            },
            "required": ["action"]
        }
    },
    {
        "name": "proxy",
        "description": "設定或取消 Windows 系統層 HTTP 代理。",
        "input_schema": {
            "type": "object",
            "properties": {
                "action": {"type": "string", "enum": ["get","set","disable"]},
                "host": {"type": "string", "description": "代理伺服器地址，例如 127.0.0.1:7890"}
            },
            "required": ["action"]
        }
    },
    {
        "name": "lock_screen",
        "description": "鎖定螢幕、登出目前使用者、切換使用者。",
        "input_schema": {
            "type": "object",
            "properties": {
                "action": {"type": "string", "enum": ["lock","logoff","switch_user"]}
            },
            "required": ["action"]
        }
    },
    {
        "name": "defender",
        "description": "Windows Defender 防毒：快速掃描、完整掃描、查看威脅歷史、新增/移除排除項目。",
        "input_schema": {
            "type": "object",
            "properties": {
                "action": {"type": "string", "enum": ["quick_scan","full_scan","status","threats","add_exclusion","remove_exclusion","list_exclusions"]},
                "path": {"type": "string", "description": "排除路徑（add/remove_exclusion 使用）"}
            },
            "required": ["action"]
        }
    },
    {
        "name": "vision_loop",
        "description": "AI 視覺自動化循環：持續截圖讓 AI 判斷畫面狀態並自動執行動作，直到達成目標或超時。真正無人值守自動化的核心。",
        "input_schema": {
            "type": "object",
            "properties": {
                "goal": {"type": "string", "description": "目標描述，例如「等待安裝完成後點擊 Finish」"},
                "max_steps": {"type": "integer", "description": "最大執行步數，預設 20"},
                "interval": {"type": "number", "description": "每步間隔秒數，預設 3"},
                "timeout": {"type": "number", "description": "總超時秒數，預設 120"}
            },
            "required": ["goal"]
        }
    },
    {
        "name": "alert_monitor",
        "description": "監控告警：持續監控 CPU/記憶體/程序/溫度/文字，條件觸發時主動發 Telegram 通知。",
        "input_schema": {
            "type": "object",
            "properties": {
                "action": {"type": "string", "enum": ["start","stop","list"]},
                "name": {"type": "string", "description": "告警名稱"},
                "condition": {"type": "string", "enum": ["cpu_above","mem_above","disk_above","process_missing","process_running","screen_text_found","temperature_above"]},
                "threshold": {"type": "string", "description": "閾值，例如 80 代表 80%"},
                "target": {"type": "string", "description": "監控目標（程序名稱或螢幕文字）"},
                "interval": {"type": "integer", "description": "檢查間隔秒數，預設 30"},
                "chat_id": {"type": "integer", "description": "通知的 Telegram chat_id，預設發給主人"}
            },
            "required": ["action"]
        }
    },
    {
        "name": "interval_schedule",
        "description": "間隔排程：每隔 N 分鐘/小時執行指定腳本或指令，可設定持續時間或執行次數。",
        "input_schema": {
            "type": "object",
            "properties": {
                "action": {"type": "string", "enum": ["start","stop","list"]},
                "name": {"type": "string", "description": "排程名稱"},
                "command": {"type": "string", "description": "要執行的指令或腳本路徑"},
                "every_minutes": {"type": "number", "description": "每隔幾分鐘執行，預設 60"},
                "repeat": {"type": "integer", "description": "執行次數，0=無限，預設 0"},
                "duration_hours": {"type": "number", "description": "持續執行小時數，0=無限，預設 0"}
            },
            "required": ["action"]
        }
    },
    {
        "name": "wait_for_text",
        "description": "等待螢幕上出現指定文字才繼續執行，常用於自動化流程中等待載入或安裝完成。",
        "input_schema": {
            "type": "object",
            "properties": {
                "text": {"type": "string", "description": "要等待的文字"},
                "timeout": {"type": "number", "description": "超時秒數，預設 60"},
                "interval": {"type": "number", "description": "檢查間隔秒數，預設 2"},
                "region": {"type": "string", "description": "監控區域 x,y,w,h（選填，預設全螢幕）"}
            },
            "required": ["text"]
        }
    },
    {
        "name": "browser_advanced",
        "description": "瀏覽器進階自動化：等待元素出現、iframe 切換、Cookie 管理、多分頁切換、表單自動填寫、下拉選單。",
        "input_schema": {
            "type": "object",
            "properties": {
                "action": {"type": "string", "enum": ["wait_element","switch_frame","get_cookies","set_cookie","list_tabs","switch_tab","new_tab","close_tab","fill_form","select_option","scroll_to","get_html","wait_url"]},
                "selector": {"type": "string", "description": "CSS 選擇器或 XPath"},
                "value": {"type": "string", "description": "填入值或 Cookie 值"},
                "name": {"type": "string", "description": "Cookie 名稱或分頁標題關鍵字"},
                "tab_index": {"type": "integer", "description": "分頁索引（switch_tab 使用）"},
                "timeout": {"type": "number", "description": "等待超時秒數，預設 30"},
                "url_pattern": {"type": "string", "description": "等待的 URL 關鍵字（wait_url 使用）"}
            },
            "required": ["action"]
        }
    },
    {
        "name": "voice_cmd",
        "description": "持續語音命令模式：背景持續監聽麥克風，說話即可控制電腦執行命令，說「停止」結束。",
        "input_schema": {
            "type": "object",
            "properties": {
                "action": {"type": "string", "enum": ["start","stop"]},
                "duration": {"type": "number", "description": "持續監聽秒數，預設 300 (5分鐘)"},
                "language": {"type": "string", "description": "語音辨識語言，預設 zh-TW"}
            },
            "required": ["action"]
        }
    },
    {
        "name": "win_notify_relay",
        "description": "Windows 通知攔截：偵測 Windows 系統通知，自動轉發到 Telegram。",
        "input_schema": {
            "type": "object",
            "properties": {
                "action": {"type": "string", "enum": ["start","stop","status"]},
                "duration": {"type": "number", "description": "監控秒數，預設 3600"},
                "filter_app": {"type": "string", "description": "只轉發指定 App 的通知（選填）"}
            },
            "required": ["action"]
        }
    },
    {
        "name": "data_process",
        "description": "JSON/CSV 資料處理：讀取、寫入、過濾、轉換、合併、統計 JSON 和 CSV 檔案。",
        "input_schema": {
            "type": "object",
            "properties": {
                "action": {"type": "string", "enum": ["read_json","write_json","read_csv","write_csv","filter","convert","merge","stats","to_table"]},
                "path": {"type": "string", "description": "檔案路徑"},
                "output": {"type": "string", "description": "輸出路徑"},
                "query": {"type": "string", "description": "過濾條件，例如 age>18"},
                "data": {"type": "string", "description": "JSON 字串資料（write 使用）"},
                "paths": {"type": "string", "description": "多個檔案路徑，逗號分隔（merge 使用）"}
            },
            "required": ["action"]
        }
    },
    {
        "name": "wake_on_lan",
        "description": "發送 WOL 魔法封包，遠端喚醒區域網路內已關機的電腦。",
        "input_schema": {
            "type": "object",
            "properties": {
                "mac": {"type": "string", "description": "目標電腦 MAC 位址，例如 AA:BB:CC:DD:EE:FF"},
                "broadcast": {"type": "string", "description": "廣播地址，預設 255.255.255.255"},
                "port": {"type": "integer", "description": "UDP 埠，預設 9"}
            },
            "required": ["mac"]
        }
    },
    {
        "name": "clipboard_history",
        "description": "剪貼簿歷史：查看多筆歷史記錄、切換到指定歷史項目、清除歷史。",
        "input_schema": {
            "type": "object",
            "properties": {
                "action": {"type": "string", "enum": ["list","get","set","clear","start_watch","stop_watch"]},
                "index": {"type": "integer", "description": "歷史索引（get/set 使用，0=最新）"}
            },
            "required": ["action"]
        }
    },
    {
        "name": "file_watcher",
        "description": "監聽資料夾檔案系統事件（新增/修改/刪除），觸發時執行命令或發送 Telegram 通知。",
        "input_schema": {
            "type": "object",
            "properties": {
                "action": {"type": "string", "enum": ["start","stop","list"]},
                "name": {"type": "string", "description": "監聽器名稱"},
                "path": {"type": "string", "description": "監聽的資料夾路徑"},
                "events": {"type": "string", "description": "監聽事件類型：created/modified/deleted/all，預設 all"},
                "command": {"type": "string", "description": "觸發時執行的指令（選填）"},
                "notify": {"type": "boolean", "description": "觸發時是否發 Telegram 通知，預設 true"}
            },
            "required": ["action"]
        }
    },
    {
        "name": "pixel_watch",
        "description": "監控螢幕指定座標的像素顏色，顏色變化時觸發動作或通知。",
        "input_schema": {
            "type": "object",
            "properties": {
                "action": {"type": "string", "enum": ["start","stop","get","list"]},
                "name": {"type": "string", "description": "監控器名稱"},
                "x": {"type": "integer"}, "y": {"type": "integer"},
                "command": {"type": "string", "description": "顏色變化時執行的指令"},
                "interval": {"type": "number", "description": "檢查間隔秒數，預設 1"},
                "tolerance": {"type": "integer", "description": "顏色容差 0-255，預設 10"}
            },
            "required": ["action"]
        }
    },
    {
        "name": "object_detect",
        "description": "AI 物件偵測：用 Claude Vision 識別螢幕上按鈕/圖示/文字區塊的精確位置，自動點擊或回傳座標。",
        "input_schema": {
            "type": "object",
            "properties": {
                "target": {"type": "string", "description": "要找的物件描述，例如「確定按鈕」「搜尋框」「關閉X」"},
                "action": {"type": "string", "enum": ["find","click","double_click"], "description": "找到後要做什麼"},
                "region": {"type": "string", "description": "搜尋區域 x,y,w,h（選填，預設全螢幕）"}
            },
            "required": ["target"]
        }
    },
    {
        "name": "mouse_record",
        "description": "滑鼠軌跡錄製與回放：錄製滑鼠移動+點擊+鍵盤完整操作，儲存後重複回放。",
        "input_schema": {
            "type": "object",
            "properties": {
                "action": {"type": "string", "enum": ["start","stop","play","list","delete"]},
                "name": {"type": "string", "description": "錄製名稱"},
                "duration": {"type": "number", "description": "錄製秒數（start 使用）"},
                "repeat": {"type": "integer", "description": "回放次數，預設 1"},
                "speed": {"type": "number", "description": "回放速度倍率，預設 1.0"}
            },
            "required": ["action"]
        }
    },
    {
        "name": "adb",
        "description": "Android 手機控制（ADB）：截圖、點擊、輸入文字、安裝 App、傳送檔案、執行 Shell 命令。",
        "input_schema": {
            "type": "object",
            "properties": {
                "action": {"type": "string", "enum": ["devices","screenshot","tap","swipe","type","key","install","push","pull","shell","start_app","stop_app"]},
                "x": {"type": "integer"}, "y": {"type": "integer"},
                "x2": {"type": "integer"}, "y2": {"type": "integer"},
                "text": {"type": "string"},
                "path": {"type": "string", "description": "本機檔案路徑"},
                "remote": {"type": "string", "description": "手機路徑"},
                "package": {"type": "string", "description": "App 套件名稱"},
                "command": {"type": "string", "description": "Shell 命令（shell 使用）"},
                "device": {"type": "string", "description": "裝置序號（多裝置時使用）"}
            },
            "required": ["action"]
        }
    },
    {
        "name": "wifi_hotspot",
        "description": "控制 Windows 行動熱點：開啟、關閉、設定 SSID 和密碼、查看連線裝置。",
        "input_schema": {
            "type": "object",
            "properties": {
                "action": {"type": "string", "enum": ["start","stop","status","set"]},
                "ssid": {"type": "string", "description": "熱點名稱（set 使用）"},
                "password": {"type": "string", "description": "熱點密碼（set 使用）"}
            },
            "required": ["action"]
        }
    },
    {
        "name": "onedrive",
        "description": "OneDrive 檔案同步：列出檔案、上傳、下載、取得分享連結。",
        "input_schema": {
            "type": "object",
            "properties": {
                "action": {"type": "string", "enum": ["list","upload","download","sync","status","open"]},
                "path": {"type": "string", "description": "本機檔案路徑"},
                "remote": {"type": "string", "description": "OneDrive 路徑（相對）"}
            },
            "required": ["action"]
        }
    },
    {
        "name": "ftp",
        "description": "FTP 客戶端：連線 FTP 伺服器，列出目錄、上傳、下載檔案。",
        "input_schema": {
            "type": "object",
            "properties": {
                "action": {"type": "string", "enum": ["list","upload","download","delete","mkdir","rename"]},
                "host": {"type": "string"},
                "user": {"type": "string"},
                "password": {"type": "string"},
                "local": {"type": "string", "description": "本機檔案路徑"},
                "remote": {"type": "string", "description": "FTP 路徑"},
                "port": {"type": "integer", "description": "FTP 埠，預設 21"}
            },
            "required": ["action","host"]
        }
    },
    {
        "name": "wsl",
        "description": "WSL（Windows Subsystem for Linux）管理：列出發行版、執行 Linux 命令、啟動/停止。",
        "input_schema": {
            "type": "object",
            "properties": {
                "action": {"type": "string", "enum": ["list","run","start","stop","status","install"]},
                "distro": {"type": "string", "description": "發行版名稱，例如 Ubuntu"},
                "command": {"type": "string", "description": "要執行的 Linux 命令（run 使用）"}
            },
            "required": ["action"]
        }
    },
    {
        "name": "hyperv",
        "description": "Hyper-V 虛擬機管理：列出、啟動、停止、暫停、建立快照、還原快照。",
        "input_schema": {
            "type": "object",
            "properties": {
                "action": {"type": "string", "enum": ["list","start","stop","pause","resume","snapshot","restore","delete_snapshot","status"]},
                "name": {"type": "string", "description": "虛擬機名稱"},
                "snapshot": {"type": "string", "description": "快照名稱（snapshot/restore 使用）"}
            },
            "required": ["action"]
        }
    },
    {
        "name": "file_diff",
        "description": "比較兩個文字檔的差異，輸出 diff 結果，可儲存為 patch 檔。",
        "input_schema": {
            "type": "object",
            "properties": {
                "file1": {"type": "string", "description": "第一個檔案路徑"},
                "file2": {"type": "string", "description": "第二個檔案路徑"},
                "output": {"type": "string", "description": "輸出 patch 檔路徑（選填）"},
                "mode": {"type": "string", "enum": ["unified","context","simple"], "description": "diff 格式，預設 unified"}
            },
            "required": ["file1","file2"]
        }
    },
    {
        "name": "screen_live",
        "description": "即時螢幕串流到 Telegram：持續截圖並發送，用於遠端監控電腦畫面。",
        "input_schema": {
            "type": "object",
            "properties": {
                "action": {"type": "string", "enum": ["start","stop"]},
                "fps": {"type": "number", "description": "每秒幾張，預設 0.5（2秒一張）"},
                "duration": {"type": "number", "description": "串流秒數，預設 60"},
                "quality": {"type": "integer", "description": "JPEG 畫質 1-95，預設 50"}
            },
            "required": ["action"]
        }
    },
    {
        "name": "download_file",
        "description": "下載網路檔案到本機指定路徑。",
        "input_schema": {
            "type": "object",
            "properties": {
                "url": {"type": "string", "description": "檔案 URL"},
                "save_path": {"type": "string", "description": "儲存路徑（選填）"}
            },
            "required": ["url"]
        }
    },
    {
        "name": "wake_listen",
        "description": "持續監聽麥克風，偵測到喚醒詞時回傳。",
        "input_schema": {
            "type": "object",
            "properties": {
                "keyword": {"type": "string", "description": "喚醒詞，預設「小牛馬」"},
                "duration": {"type": "integer", "description": "每次監聽秒數，預設 5"}
            },
            "required": []
        }
    },
    {
        "name": "right_menu",
        "description": "在指定座標右鍵點擊，並可選擇右鍵選單項目。",
        "input_schema": {
            "type": "object",
            "properties": {
                "x": {"type": "integer"},
                "y": {"type": "integer"},
                "item": {"type": "string", "description": "要選擇的選單項目文字（選填）"}
            },
            "required": ["x", "y"]
        }
    },
    {
        "name": "disk_clean",
        "description": "查看或清理 Windows 暫存資料夾，釋放磁碟空間。",
        "input_schema": {
            "type": "object",
            "properties": {
                "action": {"type": "string", "enum": ["list", "clean"], "description": "list=查看佔用，clean=清理"}
            },
            "required": ["action"]
        }
    },
    {
        "name": "usb_list",
        "description": "列出目前連接的 USB 裝置清單。",
        "input_schema": {
            "type": "object",
            "properties": {},
            "required": []
        }
    },
    {
        "name": "rdp_connect",
        "description": "使用 mstsc 連線遠端桌面（RDP）。",
        "input_schema": {
            "type": "object",
            "properties": {
                "host": {"type": "string", "description": "遠端主機 IP 或域名"},
                "user": {"type": "string", "description": "使用者名稱（選填）"},
                "width": {"type": "integer", "description": "解析度寬度，預設 1280"},
                "height": {"type": "integer", "description": "解析度高度，預設 720"}
            },
            "required": ["host"]
        }
    },
    {
        "name": "chrome_bookmarks",
        "description": "讀取並顯示 Chrome 瀏覽器的所有書籤。",
        "input_schema": {
            "type": "object",
            "properties": {},
            "required": []
        }
    },
    {
        "name": "net_share",
        "description": "網路芳鄰（SMB）：列出、連線、中斷網路磁碟機。",
        "input_schema": {
            "type": "object",
            "properties": {
                "action": {"type": "string", "enum": ["list", "connect", "disconnect"]},
                "share_path": {"type": "string", "description": "網路路徑，例如 \\\\server\\share"},
                "drive": {"type": "string", "description": "對應磁碟代號，例如 Z:"},
                "user": {"type": "string"},
                "password": {"type": "string"}
            },
            "required": ["action"]
        }
    },
    {
        "name": "font_list",
        "description": "列出系統已安裝的字型，可用關鍵字過濾。",
        "input_schema": {
            "type": "object",
            "properties": {
                "keyword": {"type": "string", "description": "搜尋關鍵字（選填）"}
            },
            "required": []
        }
    },
    {
        "name": "wait_seconds",
        "description": "等待指定秒數後繼續執行（用於自動化流程的延遲）。",
        "input_schema": {
            "type": "object",
            "properties": {
                "seconds": {"type": "number", "description": "等待秒數"}
            },
            "required": ["seconds"]
        }
    },
    {
        "name": "reminder",
        "description": "設定一次性提醒/鬧鐘，到時發出通知和語音。",
        "input_schema": {
            "type": "object",
            "properties": {
                "time": {"type": "string", "description": "HH:MM 格式時間，或純數字代表幾秒後"},
                "message": {"type": "string"}
            },
            "required": ["time", "message"]
        }
    },
    {
        "name": "webpage_shot",
        "description": "對指定網址截取整頁截圖，或監控網頁內容變化。",
        "input_schema": {
            "type": "object",
            "properties": {
                "action": {"type": "string", "enum": ["screenshot", "monitor"]},
                "url": {"type": "string"},
                "selector": {"type": "string", "description": "監控的 CSS 選擇器（monitor 使用）"},
                "interval": {"type": "number", "description": "監控間隔秒數（monitor 使用）"},
                "duration": {"type": "number", "description": "監控總秒數（monitor 使用）"}
            },
            "required": ["action", "url"]
        }
    },
    {
        "name": "file_tools",
        "description": "批次重新命名檔案、同步資料夾。",
        "input_schema": {
            "type": "object",
            "properties": {
                "action": {"type": "string", "enum": ["batch_rename", "sync"]},
                "path": {"type": "string", "description": "資料夾路徑"},
                "dest": {"type": "string", "description": "目標資料夾（sync 使用）"},
                "pattern": {"type": "string", "description": "正規表達式（batch_rename 使用）"},
                "replacement": {"type": "string", "description": "替換字串"},
                "ext": {"type": "string", "description": "副檔名過濾（batch_rename 使用）"}
            },
            "required": ["action", "path"]
        }
    },
    {
        "name": "image_tools",
        "description": "壓縮圖片、批次處理圖片資料夾、OCR+翻譯。",
        "input_schema": {
            "type": "object",
            "properties": {
                "action": {"type": "string", "enum": ["compress", "batch", "ocr_translate"]},
                "path": {"type": "string"},
                "quality": {"type": "integer", "description": "壓縮品質 0-100（compress/batch 使用）"},
                "width": {"type": "integer"}, "height": {"type": "integer"},
                "target_lang": {"type": "string", "description": "翻譯目標語言（ocr_translate 使用）"}
            },
            "required": ["action"]
        }
    },
    {
        "name": "lookup",
        "description": "查詢 IP 地理位置或外幣匯率。",
        "input_schema": {
            "type": "object",
            "properties": {
                "action": {"type": "string", "enum": ["ip", "currency"]},
                "ip": {"type": "string", "description": "IP 位址（ip 使用，留空查自己）"},
                "amount": {"type": "number", "description": "金額（currency 使用）"},
                "from_cur": {"type": "string", "description": "來源幣別，如 USD"},
                "to_cur": {"type": "string", "description": "目標幣別，如 TWD"}
            },
            "required": ["action"]
        }
    },
    {
        "name": "system_tools",
        "description": "Windows 事件日誌、USB 裝置、防火牆、印表機、網路芳鄰、字型列表。",
        "input_schema": {
            "type": "object",
            "properties": {
                "action": {"type": "string", "enum": ["event_log","usb_list","firewall_list","firewall_add","firewall_remove","printer_list","printer_jobs","net_share_list","net_share_connect","net_share_disconnect","font_list","rdp"]},
                "name": {"type": "string"},
                "host": {"type": "string"},
                "port": {"type": "integer"},
                "direction": {"type": "string", "enum": ["in","out"]},
                "share_path": {"type": "string"},
                "drive": {"type": "string"},
                "user": {"type": "string"},
                "password": {"type": "string"},
                "keyword": {"type": "string"},
                "log_name": {"type": "string"},
                "level": {"type": "string"},
                "count": {"type": "integer"}
            },
            "required": ["action"]
        }
    },
    {
        "name": "tts_advanced",
        "description": "使用微軟 Edge TTS 合成更自然的語音，支援多種中文聲音。",
        "input_schema": {
            "type": "object",
            "properties": {
                "action": {"type": "string", "enum": ["speak","list_voices"]},
                "text": {"type": "string"},
                "voice": {"type": "string", "description": "語音名稱，如 zh-CN-YunxiNeural（女）或 zh-CN-YunxiNeural（男）"}
            },
            "required": ["action"]
        }
    },
    {
        "name": "send_voice",
        "description": "將文字轉成語音並以 Telegram 語音訊息傳送給用戶。當用戶說「用語音說」、「語音回答」、「唸給我聽」、「說出來」等時使用此工具。",
        "input_schema": {
            "type": "object",
            "properties": {
                "text": {"type": "string", "description": "要轉換成語音的文字內容"},
                "voice": {"type": "string", "description": "語音名稱，預設 zh-CN-YunxiNeural（女聲）。男聲用 zh-CN-YunxiNeural"}
            },
            "required": ["text"]
        }
    },
    {
        "name": "todo_list",
        "description": "管理本地任務清單：新增、列出、完成、刪除任務。",
        "input_schema": {
            "type": "object",
            "properties": {
                "action": {"type": "string", "enum": ["add","list","done","delete","clear"]},
                "task": {"type": "string", "description": "任務內容（add 使用）"},
                "id": {"type": "integer", "description": "任務 ID（done/delete 使用）"}
            },
            "required": ["action"]
        }
    },
    {
        "name": "sysres_chart",
        "description": "即時監控 CPU/RAM 使用率並生成圖表。",
        "input_schema": {
            "type": "object",
            "properties": {
                "duration": {"type": "integer", "description": "監控秒數，預設 10"}
            },
            "required": []
        }
    },
    {
        "name": "password_mgr",
        "description": "加密儲存或查詢帳號密碼（本地加密，需主密碼）。",
        "input_schema": {
            "type": "object",
            "properties": {
                "action": {"type": "string", "enum": ["save","get"]},
                "site": {"type": "string"},
                "username": {"type": "string", "description": "帳號（save 使用）"},
                "password": {"type": "string", "description": "密碼（save 使用）"},
                "master": {"type": "string", "description": "主密碼"}
            },
            "required": ["action", "site", "master"]
        }
    },
    {
        "name": "clipboard_image",
        "description": "讀取或寫入剪貼簿中的圖片。",
        "input_schema": {
            "type": "object",
            "properties": {
                "action": {"type": "string", "enum": ["get","set"]},
                "path": {"type": "string", "description": "圖片路徑（set 使用）或儲存路徑（get 使用）"}
            },
            "required": ["action"]
        }
    },
    {
        "name": "email_control",
        "description": "讀取 IMAP 郵件收件匣。",
        "input_schema": {
            "type": "object",
            "properties": {
                "host": {"type": "string", "description": "IMAP 伺服器，如 imap.gmail.com"},
                "user": {"type": "string"},
                "password": {"type": "string"},
                "folder": {"type": "string", "description": "資料夾名稱，預設 INBOX"},
                "count": {"type": "integer", "description": "讀取封數，預設 5"}
            },
            "required": ["host", "user", "password"]
        }
    },
    {
        "name": "calendar",
        "description": "查詢或新增 Google Calendar 行事曆事件。",
        "input_schema": {
            "type": "object",
            "properties": {
                "action": {"type": "string", "enum": ["list","add","delete"]},
                "days": {"type": "integer", "description": "往後幾天（list 使用），預設 7"},
                "title": {"type": "string", "description": "行程標題（add 使用）"},
                "start": {"type": "string", "description": "開始時間，格式 2026-04-13T10:00:00（add 使用）"},
                "end": {"type": "string", "description": "結束時間（add 使用）"},
                "description": {"type": "string", "description": "說明（add 使用）"}
            },
            "required": ["action"]
        }
    },
    {
        "name": "global_hotkey",
        "description": "監聽全域快捷鍵，觸發時執行指定指令。",
        "input_schema": {
            "type": "object",
            "properties": {
                "hotkey": {"type": "string", "description": "快捷鍵組合，如 ctrl+shift+a"},
                "command": {"type": "string", "description": "觸發時執行的 shell 指令"},
                "duration": {"type": "number", "description": "監聽秒數，預設 60"}
            },
            "required": ["hotkey", "command"]
        }
    },
    {
        "name": "git",
        "description": "執行 Git 操作：狀態查看、提交、推送、拉取等。",
        "input_schema": {
            "type": "object",
            "properties": {
                "action": {"type": "string", "enum": ["status","log","pull","add","commit","push","diff"]},
                "repo": {"type": "string", "description": "repo 路徑，預設當前目錄"},
                "message": {"type": "string", "description": "commit 訊息（commit 使用）"},
                "branch": {"type": "string", "description": "分支名稱，預設 master"}
            },
            "required": ["action"]
        }
    },
    {
        "name": "hardware",
        "description": "查看進階硬體資訊：GPU 使用率/溫度、電池狀態、CPU 溫度。",
        "input_schema": {"type": "object", "properties": {}, "required": []}
    },
    {
        "name": "report",
        "description": "根據提供的資料生成 HTML 格式報告並存檔。",
        "input_schema": {
            "type": "object",
            "properties": {
                "title": {"type": "string"},
                "data": {"type": "string", "description": "JSON 格式資料，如 {\"section\": [{\"col1\": val}]}"},
                "output": {"type": "string", "description": "輸出路徑（選填）"}
            },
            "required": ["title", "data"]
        }
    },
    {
        "name": "dropbox",
        "description": "上傳或下載 Dropbox 雲端檔案。",
        "input_schema": {
            "type": "object",
            "properties": {
                "action": {"type": "string", "enum": ["upload","download"]},
                "local": {"type": "string", "description": "本地路徑"},
                "remote": {"type": "string", "description": "Dropbox 路徑，如 /folder/file.txt"},
                "token": {"type": "string", "description": "Dropbox access token（選填，可用環境變數）"}
            },
            "required": ["action", "local", "remote"]
        }
    },
    {
        "name": "docker",
        "description": "管理 Docker 容器：列出、啟動、停止、查看日誌。",
        "input_schema": {
            "type": "object",
            "properties": {
                "action": {"type": "string", "enum": ["list","start","stop","logs","images"]},
                "name": {"type": "string", "description": "容器名稱（start/stop/logs 使用）"}
            },
            "required": ["action"]
        }
    },
    {
        "name": "pdf_image",
        "description": "將 PDF 每頁轉成圖片（PNG）。",
        "input_schema": {
            "type": "object",
            "properties": {
                "path": {"type": "string"},
                "output_dir": {"type": "string", "description": "輸出資料夾（選填）"},
                "dpi": {"type": "integer", "description": "解析度，預設 150"}
            },
            "required": ["path"]
        }
    },
    {
        "name": "barcode",
        "description": "掃描圖片或截圖中的條碼、QR Code。",
        "input_schema": {
            "type": "object",
            "properties": {
                "image_path": {"type": "string", "description": "圖片路徑（選填，不填則截圖）"}
            },
            "required": []
        }
    },
    {
        "name": "nlp",
        "description": "AI 文字分析：摘要長文或分析情緒。",
        "input_schema": {
            "type": "object",
            "properties": {
                "action": {"type": "string", "enum": ["summarize","sentiment"]},
                "text": {"type": "string"}
            },
            "required": ["action", "text"]
        }
    },
    {
        "name": "vpn",
        "description": "VPN 連線管理：列出、連線、斷線。",
        "input_schema": {
            "type": "object",
            "properties": {
                "action": {"type": "string", "enum": ["list","connect","disconnect"]},
                "name": {"type": "string", "description": "VPN 名稱"},
                "user": {"type": "string"},
                "password": {"type": "string"}
            },
            "required": ["action"]
        }
    },
    {
        "name": "restore_point",
        "description": "建立或列出 Windows 系統還原點。",
        "input_schema": {
            "type": "object",
            "properties": {
                "action": {"type": "string", "enum": ["create","list"]},
                "description": {"type": "string", "description": "還原點說明（create 使用）"}
            },
            "required": ["action"]
        }
    },
    {
        "name": "disk_analyze",
        "description": "分析磁碟空間使用，列出佔用最多的資料夾。",
        "input_schema": {
            "type": "object",
            "properties": {
                "path": {"type": "string", "description": "分析路徑，預設 C:/"},
                "top": {"type": "integer", "description": "顯示前幾名，預設 10"}
            },
            "required": []
        }
    },
    {
        "name": "face_detect",
        "description": "從截圖或圖片中偵測人臉，並標記位置。",
        "input_schema": {
            "type": "object",
            "properties": {
                "image_path": {"type": "string", "description": "圖片路徑（選填，不填則截圖）"},
                "output": {"type": "string", "description": "輸出路徑（選填）"}
            },
            "required": []
        }
    },
    {
        "name": "video_gif",
        "description": "將影片片段轉成 GIF 動圖。",
        "input_schema": {
            "type": "object",
            "properties": {
                "path": {"type": "string"},
                "start": {"type": "number", "description": "起始秒，預設 0"},
                "duration": {"type": "number", "description": "持續秒數，預設 5"},
                "output": {"type": "string", "description": "輸出路徑（選填）"},
                "fps": {"type": "integer", "description": "GIF fps，預設 10"}
            },
            "required": ["path"]
        }
    },
    {
        "name": "excel_chart",
        "description": "在 Excel 工作表中生成圖表（長條/折線/圓餅）。",
        "input_schema": {
            "type": "object",
            "properties": {
                "path": {"type": "string"},
                "sheet": {"type": "string"},
                "type": {"type": "string", "enum": ["bar","line","pie"], "description": "圖表類型，預設 bar"},
                "title": {"type": "string"}
            },
            "required": ["path", "sheet"]
        }
    },
    {
        "name": "speedtest",
        "description": "測試目前網路的上下載速度和 Ping 值。",
        "input_schema": {"type": "object", "properties": {}, "required": []}
    },
    {
        "name": "screenshot_compare",
        "description": "比對兩張截圖或圖片的差異，標記變化區域。",
        "input_schema": {
            "type": "object",
            "properties": {
                "img1": {"type": "string", "description": "第一張圖片路徑（選填，不填則即時截圖）"},
                "img2": {"type": "string", "description": "第二張圖片路徑（選填，不填則 2 秒後再截圖）"},
                "output": {"type": "string", "description": "輸出路徑（選填）"}
            },
            "required": []
        }
    },
    {
        "name": "screen_record",
        "description": "螢幕錄影或攝影機拍照。",
        "input_schema": {
            "type": "object",
            "properties": {
                "action": {"type": "string", "enum": ["record","webcam"]},
                "duration": {"type": "number", "description": "錄影秒數（record 使用）"},
                "output": {"type": "string", "description": "輸出檔案路徑（選填）"}
            },
            "required": ["action"]
        }
    },
    {
        "name": "translate",
        "description": "翻譯文字，支援中英日韓等多語言。",
        "input_schema": {
            "type": "object",
            "properties": {
                "text": {"type": "string"},
                "target": {"type": "string", "description": "目標語言代碼，如 zh-TW, en, ja, ko，預設 zh-TW"},
                "source": {"type": "string", "description": "來源語言，預設 auto 自動偵測"}
            },
            "required": ["text"]
        }
    },
    {
        "name": "chart",
        "description": "根據數據生成折線圖、長條圖、圓餅圖並存成圖片。",
        "input_schema": {
            "type": "object",
            "properties": {
                "type": {"type": "string", "enum": ["line","bar","pie"]},
                "data": {"type": "string", "description": "JSON 格式資料，如 {\"A\": [1,2,3]} 或 {\"A\": 30, \"B\": 70}"},
                "title": {"type": "string", "description": "圖表標題（選填）"},
                "output": {"type": "string", "description": "輸出路徑（選填）"}
            },
            "required": ["type", "data"]
        }
    },
    {
        "name": "pptx_control",
        "description": "讀取或建立 PowerPoint 簡報。",
        "input_schema": {
            "type": "object",
            "properties": {
                "action": {"type": "string", "enum": ["read","create"]},
                "path": {"type": "string"},
                "slides": {"type": "string", "description": "JSON 格式投影片，如 [{\"title\":\"標題\",\"body\":\"內容\"}]（create 使用）"}
            },
            "required": ["action", "path"]
        }
    },
    {
        "name": "api_call",
        "description": "呼叫任意外部 HTTP REST API。",
        "input_schema": {
            "type": "object",
            "properties": {
                "method": {"type": "string", "enum": ["GET","POST","PUT","DELETE","PATCH"]},
                "url": {"type": "string"},
                "headers": {"type": "string", "description": "JSON 格式 headers（選填）"},
                "body": {"type": "string", "description": "JSON 格式 body（選填）"}
            },
            "required": ["method", "url"]
        }
    },
    {
        "name": "watchdog",
        "description": "守護指定程序，若崩潰則自動重啟。",
        "input_schema": {
            "type": "object",
            "properties": {
                "process": {"type": "string", "description": "程序名稱，如 pythonw.exe"},
                "script": {"type": "string", "description": "崩潰時要執行的腳本路徑"},
                "duration": {"type": "number", "description": "守護秒數，預設 60"}
            },
            "required": ["process", "script"]
        }
    },
    {
        "name": "ssh_sftp",
        "description": "SSH 遠端執行指令，或 SFTP 上傳/下載檔案。",
        "input_schema": {
            "type": "object",
            "properties": {
                "action": {"type": "string", "enum": ["ssh_run","sftp_upload","sftp_download"]},
                "host": {"type": "string"},
                "user": {"type": "string"},
                "password": {"type": "string"},
                "command": {"type": "string", "description": "SSH 指令（ssh_run 使用）"},
                "local": {"type": "string", "description": "本地路徑"},
                "remote": {"type": "string", "description": "遠端路徑"}
            },
            "required": ["action", "host", "user", "password"]
        }
    },
    {
        "name": "network_diag",
        "description": "網路診斷：Ping、路由追蹤、Port 掃描。",
        "input_schema": {
            "type": "object",
            "properties": {
                "action": {"type": "string", "enum": ["ping","traceroute","portscan"]},
                "host": {"type": "string"},
                "ports": {"type": "string", "description": "要掃描的 port 列表，如 22,80,443（portscan 使用）"}
            },
            "required": ["action", "host"]
        }
    },
    {
        "name": "win_service",
        "description": "管理 Windows 系統服務：列出、啟動、停止。",
        "input_schema": {
            "type": "object",
            "properties": {
                "action": {"type": "string", "enum": ["list","start","stop"]},
                "name": {"type": "string", "description": "服務名稱（start/stop 使用）"}
            },
            "required": ["action"]
        }
    },
    {
        "name": "pdf_edit",
        "description": "PDF 進階編輯：合併、分割、加浮水印。",
        "input_schema": {
            "type": "object",
            "properties": {
                "action": {"type": "string", "enum": ["merge","split","watermark"]},
                "path": {"type": "string", "description": "PDF 路徑"},
                "output": {"type": "string", "description": "輸出路徑"},
                "paths": {"type": "string", "description": "JSON 格式多個 PDF 路徑（merge 使用）"},
                "text": {"type": "string", "description": "浮水印文字（watermark 使用）"}
            },
            "required": ["action"]
        }
    },
    {
        "name": "audio_process",
        "description": "音訊處理：格式轉換、剪輯。",
        "input_schema": {
            "type": "object",
            "properties": {
                "action": {"type": "string", "enum": ["convert","trim"]},
                "input": {"type": "string", "description": "輸入檔案路徑"},
                "output": {"type": "string", "description": "輸出檔案路徑"},
                "start_ms": {"type": "integer", "description": "起始毫秒（trim 使用）"},
                "end_ms": {"type": "integer", "description": "結束毫秒（trim 使用）"}
            },
            "required": ["action", "input"]
        }
    },
    {
        "name": "push_notify",
        "description": "發送推播通知到 Discord 或 LINE。",
        "input_schema": {
            "type": "object",
            "properties": {
                "platform": {"type": "string", "enum": ["discord","line"]},
                "message": {"type": "string"},
                "webhook_or_token": {"type": "string", "description": "Discord Webhook URL 或 LINE Notify Token"}
            },
            "required": ["platform", "message", "webhook_or_token"]
        }
    },
    {
        "name": "disk_backup",
        "description": "磁碟暫存清理或資料夾備份。",
        "input_schema": {
            "type": "object",
            "properties": {
                "action": {"type": "string", "enum": ["list_temp","clean_temp","backup"]},
                "src": {"type": "string", "description": "備份來源路徑（backup 使用）"},
                "dest": {"type": "string", "description": "備份目標資料夾（backup 使用）"}
            },
            "required": ["action"]
        }
    },
    {
        "name": "registry",
        "description": "讀取或寫入 Windows 登錄檔。",
        "input_schema": {
            "type": "object",
            "properties": {
                "action": {"type": "string", "enum": ["read","write"]},
                "key": {"type": "string", "description": "登錄檔路徑，如 HKCU\\Software\\MyApp"},
                "value_name": {"type": "string", "description": "值名稱"},
                "value": {"type": "string", "description": "要寫入的值（write 使用）"}
            },
            "required": ["action", "key"]
        }
    },
    {
        "name": "ai_video",
        "description": "用 AI API 生成影片（Replicate / Runway / Kling）。需要對應的 API Key 在 .env 中。",
        "input_schema": {
            "type": "object",
            "properties": {
                "prompt": {"type": "string", "description": "影片描述文字"},
                "provider": {
                    "type": "string",
                    "enum": ["replicate", "runway", "kling"],
                    "description": "AI 服務商（replicate=最多模型, runway=高品質, kling=快影）"
                },
                "model": {"type": "string", "description": "模型名稱（replicate 用，如 minimax/video-01）"},
                "image_url": {"type": "string", "description": "起始圖片 URL（image-to-video，選填）"},
                "duration": {"type": "number", "description": "影片秒數（預設 5）"},
                "output": {"type": "string", "description": "輸出路徑（選填，預設桌面）"}
            },
            "required": ["prompt", "provider"]
        }
    },
    {
        "name": "video_gen",
        "description": "生成影片：投影片、文字動畫、TTS語音影片、螢幕錄影。",
        "input_schema": {
            "type": "object",
            "properties": {
                "mode": {
                    "type": "string",
                    "enum": ["slideshow", "text_video", "tts_video", "screen_record"],
                    "description": "slideshow=圖片投影片, text_video=文字動畫, tts_video=語音影片, screen_record=螢幕錄影"
                },
                "output": {"type": "string", "description": "輸出路徑（選填，預設桌面）"},
                "text": {"type": "string", "description": "文字內容（text_video/tts_video使用）"},
                "images": {"type": "array", "items": {"type": "string"}, "description": "圖片路徑列表（slideshow使用）"},
                "image": {"type": "string", "description": "背景圖片路徑（tts_video使用，選填）"},
                "duration": {"type": "number", "description": "時長秒數（text_video/screen_record）或每張停留秒數（slideshow）"},
                "fps": {"type": "number", "description": "幀率（選填，預設24）"},
                "voice": {"type": "string", "description": "TTS聲音（tts_video，預設zh-CN-YunxiNeural）"},
                "bg_color": {"type": "array", "description": "背景顏色RGB（text_video）"},
                "font_color": {"type": "array", "description": "字體顏色RGB（text_video）"},
                "font_size": {"type": "number", "description": "字體大小（text_video）"},
                "subtitle": {"type": "boolean", "description": "是否顯示字幕（tts_video，預設true）"}
            },
            "required": ["mode"]
        }
    },
    {
        "name": "video_process",
        "description": "影片處理：截取指定秒數畫面、剪輯片段。",
        "input_schema": {
            "type": "object",
            "properties": {
                "action": {"type": "string", "enum": ["screenshot","trim"]},
                "path": {"type": "string"},
                "second": {"type": "number", "description": "截取的秒數（screenshot 使用）"},
                "start": {"type": "number", "description": "起始秒（trim 使用）"},
                "end": {"type": "number", "description": "結束秒（trim 使用）"},
                "output": {"type": "string", "description": "輸出路徑（選填）"}
            },
            "required": ["action", "path"]
        }
    },
    {
        "name": "monitor_config",
        "description": "列出所有螢幕顯示器的資訊（解析度、位置、主副螢幕）。",
        "input_schema": {"type": "object", "properties": {}, "required": []}
    },
    {
        "name": "run_code",
        "description": "直接執行 Python 程式碼或 PowerShell 指令。",
        "input_schema": {
            "type": "object",
            "properties": {
                "type": {"type": "string", "enum": ["python", "shell"], "description": "python=執行Python, shell=執行PowerShell"},
                "code": {"type": "string", "description": "要執行的程式碼或指令"}
            },
            "required": ["type", "code"]
        }
    },
    {
        "name": "document_control",
        "description": "讀取或寫入 Word、Excel、PDF 文件。",
        "input_schema": {
            "type": "object",
            "properties": {
                "action": {"type": "string", "enum": ["word_read","word_write","excel_read","excel_write","pdf_read"]},
                "path": {"type": "string", "description": "文件路徑"},
                "content": {"type": "string", "description": "要寫入的內容（word_write/excel_write 使用）"},
                "sheet": {"type": "string", "description": "Excel 工作表名稱（excel_read/excel_write 使用）"}
            },
            "required": ["action", "path"]
        }
    },
    {
        "name": "web_scrape",
        "description": "讀取/瀏覽網頁內容並回傳文字。當用戶要求「去看看這個網址」「瀏覽網頁」「讀取網站內容」「網頁裡面有什麼」時，必須使用此工具（action=scrape），不要使用 browser_control。支援 JS 動態網站。",
        "input_schema": {
            "type": "object",
            "properties": {
                "action": {"type": "string", "enum": ["scrape", "screen_diff"]},
                "url": {"type": "string", "description": "網址（scrape 使用）"},
                "selector": {"type": "string", "description": "CSS 選擇器（scrape 使用）"},
                "interval": {"type": "number", "description": "偵測間隔秒數（screen_diff 使用）"},
                "region": {"type": "string", "description": "螢幕區域（screen_diff 使用）"}
            },
            "required": ["action"]
        }
    },
    {
        "name": "image_edit",
        "description": "圖片編輯：裁切、縮放、加文字、合併圖片。",
        "input_schema": {
            "type": "object",
            "properties": {
                "action": {"type": "string", "enum": ["crop","resize","text","merge"]},
                "path": {"type": "string", "description": "圖片路徑"},
                "params": {"type": "string", "description": "參數（crop: x y w h; resize: w h; text: x y 文字; merge: 圖片2路徑）"}
            },
            "required": ["action", "path"]
        }
    },
    {
        "name": "cloud_storage",
        "description": "上傳或下載 Google Drive 檔案。",
        "input_schema": {
            "type": "object",
            "properties": {
                "action": {"type": "string", "enum": ["upload","download"]},
                "path": {"type": "string", "description": "本地檔案路徑"},
                "drive_id": {"type": "string", "description": "Google Drive 檔案或資料夾 ID"}
            },
            "required": ["action", "path"]
        }
    },
    {
        "name": "database",
        "description": "執行 SQLite 或 MySQL 資料庫查詢。",
        "input_schema": {
            "type": "object",
            "properties": {
                "type": {"type": "string", "enum": ["sqlite","mysql"]},
                "db": {"type": "string", "description": "SQLite: 資料庫路徑；MySQL: host"},
                "name": {"type": "string", "description": "MySQL 資料庫名稱"},
                "sql": {"type": "string", "description": "SQL 指令"}
            },
            "required": ["type", "db", "sql"]
        }
    },
    {
        "name": "encrypt_file",
        "description": "加密或解密檔案。",
        "input_schema": {
            "type": "object",
            "properties": {
                "action": {"type": "string", "enum": ["encrypt","decrypt"]},
                "path": {"type": "string", "description": "檔案路徑"},
                "password": {"type": "string", "description": "密碼"}
            },
            "required": ["action", "path", "password"]
        }
    },
    {
        "name": "qr_code",
        "description": "生成或掃描 QR Code，或監控剪貼簿變化。",
        "input_schema": {
            "type": "object",
            "properties": {
                "action": {"type": "string", "enum": ["qr_gen","qr_scan","clipboard_watch"]},
                "content": {"type": "string", "description": "QR Code 內容（qr_gen 使用）"},
                "path": {"type": "string", "description": "圖片路徑（qr_scan 使用）或儲存路徑（qr_gen 使用）"},
                "duration": {"type": "number", "description": "監控秒數（clipboard_watch 使用）"}
            },
            "required": ["action"]
        }
    },
    # ── 缺口1：觸發驅動 ───────────────────────────────────────────
    {
        "name": "email_trigger",
        "description": "監控 Email 收件箱，當收到符合條件的郵件時自動回傳內容或觸發動作。支援 Gmail/IMAP。",
        "input_schema": {
            "type": "object",
            "properties": {
                "action": {"type": "string", "enum": ["check","watch","send"], "description": "check=立即檢查, watch=持續監控N秒, send=發送郵件"},
                "host": {"type": "string", "description": "IMAP 伺服器，如 imap.gmail.com"},
                "user": {"type": "string", "description": "Email 帳號"},
                "password": {"type": "string", "description": "Email 密碼或應用程式密碼"},
                "filter_from": {"type": "string", "description": "過濾寄件人（選填）"},
                "filter_subject": {"type": "string", "description": "過濾主旨關鍵字（選填）"},
                "duration": {"type": "number", "description": "監控秒數（watch 使用，預設 300）"},
                "to": {"type": "string", "description": "收件人（send 使用）"},
                "subject": {"type": "string", "description": "郵件主旨（send 使用）"},
                "body": {"type": "string", "description": "郵件內容（send 使用）"}
            },
            "required": ["action"]
        }
    },
    {
        "name": "file_trigger",
        "description": "監控資料夾，當有新增/修改/刪除檔案時自動執行指定動作（複製、通知、執行程式）。",
        "input_schema": {
            "type": "object",
            "properties": {
                "folder": {"type": "string", "description": "監控的資料夾路徑"},
                "event": {"type": "string", "enum": ["created","modified","deleted","any"], "description": "觸發事件類型"},
                "pattern": {"type": "string", "description": "檔案名稱 glob 過濾，如 *.pdf（選填）"},
                "action": {"type": "string", "enum": ["notify","copy","run","list"], "description": "觸發後動作"},
                "target": {"type": "string", "description": "複製目標資料夾 或 要執行的程式路徑（選填）"},
                "duration": {"type": "number", "description": "監控秒數（預設 60）"}
            },
            "required": ["folder","event","action"]
        }
    },
    {
        "name": "webhook_server",
        "description": "在本機啟動或停止 Webhook HTTP 伺服器，外部服務可推送事件觸發動作；也可作為遠端手機控制入口。",
        "input_schema": {
            "type": "object",
            "properties": {
                "action": {"type": "string", "enum": ["start","stop","status","test"], "description": "start=啟動, stop=停止, status=狀態, test=發測試請求"},
                "port": {"type": "number", "description": "監聽 port（預設 8765）"},
                "secret": {"type": "string", "description": "驗證 token（選填）"}
            },
            "required": ["action"]
        }
    },
    # ── 缺口2：應用程式深度控制 ────────────────────────────────────
    {
        "name": "com_auto",
        "description": "透過 Windows COM 介面深度控制 Excel/Word/Outlook，執行巨集、讀寫儲存格、寄信、建立文件等。",
        "input_schema": {
            "type": "object",
            "properties": {
                "app": {"type": "string", "enum": ["excel","word","outlook","powerpoint"], "description": "目標應用程式"},
                "action": {"type": "string", "description": "excel: read_cell/write_cell/run_macro/save/open/close/list_sheets\nword: read/write/save/open/close\noutlook: send/list_inbox/read_mail"},
                "path": {"type": "string", "description": "檔案路徑（選填）"},
                "sheet": {"type": "string", "description": "工作表名稱（Excel 選填）"},
                "cell": {"type": "string", "description": "儲存格，如 A1（Excel 選填）"},
                "value": {"type": "string", "description": "寫入值或郵件內容（選填）"},
                "macro": {"type": "string", "description": "巨集名稱（run_macro 使用）"},
                "to": {"type": "string", "description": "收件人（Outlook send 使用）"},
                "subject": {"type": "string", "description": "主旨（Outlook send 使用）"}
            },
            "required": ["app","action"]
        }
    },
    {
        "name": "dialog_auto",
        "description": "自動偵測並處理系統對話框（UAC、MessageBox、確認視窗），自動點擊指定按鈕。",
        "input_schema": {
            "type": "object",
            "properties": {
                "action": {"type": "string", "enum": ["find_and_click","list_dialogs","wait_and_click"], "description": "find_and_click=找到並點擊, list_dialogs=列出所有對話框, wait_and_click=等待出現後點擊"},
                "button_text": {"type": "string", "description": "要點擊的按鈕文字，如 確定/OK/是/Yes（選填，預設自動選確認）"},
                "window_title": {"type": "string", "description": "目標視窗標題關鍵字（選填）"},
                "timeout": {"type": "number", "description": "等待秒數（wait_and_click 使用，預設 30）"}
            },
            "required": ["action"]
        }
    },
    {
        "name": "ime_switch",
        "description": "切換輸入法（中文/英文），或查詢目前輸入法狀態。",
        "input_schema": {
            "type": "object",
            "properties": {
                "action": {"type": "string", "enum": ["switch_en","switch_zh","toggle","status"], "description": "switch_en=切英文, switch_zh=切中文, toggle=切換, status=查詢目前狀態"}
            },
            "required": ["action"]
        }
    },
    # ── 缺口3：感知能力 ────────────────────────────────────────────
    {
        "name": "wake_word",
        "description": "持續監聽麥克風，偵測到喚醒詞後回傳觸發事件，或持續錄音轉文字。",
        "input_schema": {
            "type": "object",
            "properties": {
                "action": {"type": "string", "enum": ["listen_once","transcribe_stream","detect_keyword"], "description": "listen_once=錄一句話轉文字, transcribe_stream=持續轉錄N秒, detect_keyword=偵測關鍵字觸發"},
                "keyword": {"type": "string", "description": "偵測關鍵字（detect_keyword 使用）"},
                "duration": {"type": "number", "description": "錄音/監聽秒數（預設 5）"},
                "language": {"type": "string", "description": "語言代碼，如 zh-TW/en-US（預設 zh-TW）"}
            },
            "required": ["action"]
        }
    },
    {
        "name": "sound_detect",
        "description": "分析麥克風或系統音訊，偵測音量超過閾值、特定頻率或環境音類型（靜音/說話/音樂）。",
        "input_schema": {
            "type": "object",
            "properties": {
                "action": {"type": "string", "enum": ["volume_level","detect_silence","detect_speech","record_until_silence"], "description": "volume_level=取得當前音量, detect_silence=偵測靜音, detect_speech=偵測說話, record_until_silence=錄音直到靜音"},
                "threshold": {"type": "number", "description": "音量閾值 0-100（選填，預設 20）"},
                "duration": {"type": "number", "description": "監聽秒數（選填，預設 5）"},
                "output": {"type": "string", "description": "錄音儲存路徑（選填）"}
            },
            "required": ["action"]
        }
    },
    {
        "name": "face_recognize",
        "description": "用攝影機拍照並辨識人臉，可訓練登記人臉或驗證是否為指定人物。",
        "input_schema": {
            "type": "object",
            "properties": {
                "action": {"type": "string", "enum": ["detect","recognize","enroll","capture"], "description": "detect=偵測人臉數量, recognize=識別是誰, enroll=登記新人臉, capture=拍照儲存"},
                "name": {"type": "string", "description": "人物名稱（enroll/recognize 使用）"},
                "image_path": {"type": "string", "description": "使用現有圖片而非攝影機（選填）"},
                "output": {"type": "string", "description": "儲存路徑（capture/enroll 使用，選填）"}
            },
            "required": ["action"]
        }
    },
    # ── 缺口4：跨裝置控制 ──────────────────────────────────────────
    {
        "name": "http_server",
        "description": "在本機啟動 HTTP 控制伺服器，手機/其他裝置可透過瀏覽器或 API 遠端發送指令控制電腦。",
        "input_schema": {
            "type": "object",
            "properties": {
                "action": {"type": "string", "enum": ["start","stop","status","get_url"], "description": "start=啟動, stop=停止, status=狀態, get_url=取得存取網址"},
                "port": {"type": "number", "description": "監聽 port（預設 9876）"},
                "password": {"type": "string", "description": "存取密碼（選填）"}
            },
            "required": ["action"]
        }
    },
    {
        "name": "lan_scan",
        "description": "掃描本機區域網路，列出連線的裝置 IP、MAC 位址、主機名稱，或測試特定 port 是否開啟。",
        "input_schema": {
            "type": "object",
            "properties": {
                "action": {"type": "string", "enum": ["scan","ping_sweep","port_check","get_local_ip"], "description": "scan=掃描區網, ping_sweep=ping 掃描, port_check=檢查特定 port, get_local_ip=取得本機 IP"},
                "subnet": {"type": "string", "description": "要掃描的子網路，如 192.168.1.0/24（選填，自動偵測）"},
                "host": {"type": "string", "description": "目標主機（port_check 使用）"},
                "port": {"type": "number", "description": "要檢查的 port（port_check 使用）"}
            },
            "required": ["action"]
        }
    },
    {
        "name": "serial_port",
        "description": "透過 Serial/COM Port 與 Arduino、IoT 裝置、感測器通訊，發送/接收串列資料。",
        "input_schema": {
            "type": "object",
            "properties": {
                "action": {"type": "string", "enum": ["list","open","close","send","read","send_read"], "description": "list=列出 COM port, send=發送資料, read=讀取資料, send_read=發送後讀取回應"},
                "port": {"type": "string", "description": "COM port 名稱，如 COM3（除 list 外必填）"},
                "baudrate": {"type": "number", "description": "鮑率（預設 9600）"},
                "data": {"type": "string", "description": "要發送的資料（send/send_read 使用）"},
                "timeout": {"type": "number", "description": "讀取逾時秒數（預設 2）"}
            },
            "required": ["action"]
        }
    },
    {
        "name": "mqtt",
        "description": "MQTT 訊息發布/訂閱，用於智慧家庭、IoT 裝置控制（Node-RED、Home Assistant 等）。",
        "input_schema": {
            "type": "object",
            "properties": {
                "action": {"type": "string", "enum": ["publish","subscribe","test_connect"], "description": "publish=發布訊息, subscribe=訂閱主題, test_connect=測試連線"},
                "broker": {"type": "string", "description": "MQTT broker 位址，如 192.168.1.1"},
                "port": {"type": "number", "description": "broker port（預設 1883）"},
                "topic": {"type": "string", "description": "MQTT 主題"},
                "message": {"type": "string", "description": "要發布的訊息（publish 使用）"},
                "duration": {"type": "number", "description": "訂閱持續秒數（subscribe 使用，預設 10）"},
                "username": {"type": "string", "description": "帳號（選填）"},
                "password": {"type": "string", "description": "密碼（選填）"}
            },
            "required": ["action","broker"]
        }
    },
    # ── 缺口5：內容理解與處理 ──────────────────────────────────────
    {
        "name": "doc_ai",
        "description": "用 AI 深度理解文件/圖片，自動提取發票金額、合約條款、報表數據、名片資訊等結構化資料。",
        "input_schema": {
            "type": "object",
            "properties": {
                "action": {"type": "string", "enum": ["extract","summarize","classify","compare","qa"], "description": "extract=提取欄位, summarize=摘要, classify=分類, compare=比較兩份文件, qa=問答"},
                "path": {"type": "string", "description": "文件或圖片路徑（PDF/Word/圖片）"},
                "path2": {"type": "string", "description": "第二份文件（compare 使用）"},
                "fields": {"type": "string", "description": "要提取的欄位，逗號分隔，如 金額,日期,供應商（extract 使用）"},
                "question": {"type": "string", "description": "對文件提問（qa 使用）"},
                "url": {"type": "string", "description": "文件 URL（選填，可代替 path）"}
            },
            "required": ["action"]
        }
    },
    {
        "name": "web_monitor",
        "description": "定期監控網頁內容變化，偵測到變化時通知（適合監控商品價格、職缺、新聞、股票公告）。",
        "input_schema": {
            "type": "object",
            "properties": {
                "action": {"type": "string", "enum": ["check_once","watch","diff","get_price"], "description": "check_once=立即抓取內容, watch=持續監控N秒, diff=比較與上次差異, get_price=提取價格數字"},
                "url": {"type": "string", "description": "要監控的網頁 URL"},
                "selector": {"type": "string", "description": "CSS selector 選取特定區塊（選填，預設 body）"},
                "interval": {"type": "number", "description": "檢查間隔秒數（watch 使用，預設 60）"},
                "duration": {"type": "number", "description": "監控總時長秒數（watch 使用，預設 300）"},
                "keyword": {"type": "string", "description": "偵測特定關鍵字出現/消失（選填）"}
            },
            "required": ["action","url"]
        }
    },
    {
        "name": "audio_transcribe",
        "description": "即時轉錄系統播放的音訊或麥克風輸入，也可轉錄已有的音訊檔案（支援中英文）。",
        "input_schema": {
            "type": "object",
            "properties": {
                "action": {"type": "string", "enum": ["transcribe_file","transcribe_mic","transcribe_system"], "description": "transcribe_file=轉錄音訊檔, transcribe_mic=麥克風錄音轉錄, transcribe_system=轉錄系統播放音訊"},
                "path": {"type": "string", "description": "音訊檔路徑（transcribe_file 使用）"},
                "duration": {"type": "number", "description": "錄製秒數（transcribe_mic/system 使用，預設 30）"},
                "language": {"type": "string", "description": "語言，如 zh/en（預設自動偵測）"},
                "output": {"type": "string", "description": "轉錄文字儲存路徑（選填）"}
            },
            "required": ["action"]
        }
    },
    # ══════════════════════════════════════════════════
    # 奧創升級技能集 TOOLS
    # ══════════════════════════════════════════════════
    {
        "name": "osint_search",
        "description": "OSINT情報蒐集。搜尋網路資訊、新聞、Reddit社群、查詢IP/域名情報、取得頭條新聞。當用戶要查某人/公司/IP/域名的情報，或想搜尋特定資訊時使用。",
        "input_schema": {
            "type": "object",
            "properties": {
                "action": {"type": "string", "enum": ["web_search","news_search","reddit_search","ip_osint","domain_osint","top_news"],
                           "description": "web_search=網頁搜尋, news_search=新聞搜尋, reddit_search=Reddit搜尋, ip_osint=IP情報, domain_osint=域名情報, top_news=台灣頭條"},
                "query": {"type": "string", "description": "搜尋關鍵字（web/news/reddit_search 使用）"},
                "target": {"type": "string", "description": "目標IP或域名（ip_osint/domain_osint 使用）"},
                "limit": {"type": "number", "description": "結果數量限制（預設10）"}
            },
            "required": ["action"]
        }
    },
    {
        "name": "news_monitor",
        "description": "新聞監控。查詢最新頭條、搜尋特定關鍵字新聞、開始背景監控新關鍵字新聞。當用戶問新聞/時事/最新消息時使用。",
        "input_schema": {
            "type": "object",
            "properties": {
                "action": {"type": "string", "enum": ["check","top_headlines","start_watch"],
                           "description": "check=查詢關鍵字新聞, top_headlines=頭條新聞, start_watch=開始背景監控"},
                "keywords": {"type": "string", "description": "關鍵字，多個用逗號分隔（如 台積電,AI,台灣）"},
                "interval": {"type": "number", "description": "監控間隔秒數（start_watch 使用，預設 300）"},
                "duration": {"type": "number", "description": "監控總時長秒數（start_watch 使用，預設 3600）"}
            },
            "required": ["action"]
        }
    },
    {
        "name": "threat_intel",
        "description": "威脅情報分析。透過 VirusTotal 分析URL/IP/Hash是否惡意，查詢IP濫用記錄，掃描目前電腦的外部連線。用於資安分析。",
        "input_schema": {
            "type": "object",
            "properties": {
                "action": {"type": "string", "enum": ["check_url","check_ip","check_hash","check_abuse_ip","scan_connections"],
                           "description": "check_url=分析URL, check_ip=分析IP, check_hash=分析檔案Hash, check_abuse_ip=IP濫用查詢, scan_connections=掃描當前外部連線"},
                "target": {"type": "string", "description": "分析目標（URL/IP/Hash）"},
                "api_key": {"type": "string", "description": "VirusTotal/AbuseIPDB API Key（選填，優先用環境變數）"}
            },
            "required": ["action"]
        }
    },
    {
        "name": "auto_skill",
        "description": "自動技能生成與部署。讓AI自動寫新程式技能、測試、部署到bot。列出現有技能清單。當用戶要新增功能或查詢bot有哪些技能時使用。",
        "input_schema": {
            "type": "object",
            "properties": {
                "action": {"type": "string", "enum": ["generate","test","deploy","list_skills"],
                           "description": "generate=AI生成技能, test=測試程式碼, deploy=部署技能, list_skills=列出所有技能"},
                "goal": {"type": "string", "description": "技能功能需求描述（generate 使用）"},
                "skill_name": {"type": "string", "description": "技能英文名稱（如 send_sms）"},
                "code": {"type": "string", "description": "要測試或部署的程式碼（test/deploy 使用）"}
            },
            "required": ["action"]
        }
    },
    {
        "name": "smart_home",
        "description": "智慧家居控制。透過 Home Assistant 控制燈光、插座、空調、場景等所有智慧家電。當用戶要控制家電、查詢設備狀態時使用。",
        "input_schema": {
            "type": "object",
            "properties": {
                "action": {"type": "string", "enum": ["list_devices","turn_on","turn_off","get_state","set_value","run_scene"],
                           "description": "list_devices=列出所有設備, turn_on=開啟, turn_off=關閉, get_state=查狀態, set_value=設定數值, run_scene=執行場景"},
                "device": {"type": "string", "description": "設備entity_id（如 light.living_room, switch.fan）或場景名稱"},
                "value": {"type": "string", "description": "設定值（set_value 使用）"},
                "host": {"type": "string", "description": "HA主機地址（選填，預設用 HA_HOST 環境變數）"},
                "token": {"type": "string", "description": "HA長效Token（選填，預設用 HA_TOKEN 環境變數）"}
            },
            "required": ["action"]
        }
    },
    {
        "name": "goal_manager",
        "description": "目標管理系統。設定長期目標、AI自動拆解執行步驟、追蹤進度、查看下一個待執行任務。當用戶有計畫要實現、想追蹤目標時使用。",
        "input_schema": {
            "type": "object",
            "properties": {
                "action": {"type": "string", "enum": ["add","list","detail","update_status","set_progress","delete","next_task"],
                           "description": "add=新增目標, list=列出所有, detail=查看詳情, update_status=更新狀態, set_progress=設定進度, delete=刪除, next_task=下一個待辦"},
                "goal": {"type": "string", "description": "目標描述（add 使用）"},
                "goal_id": {"type": "string", "description": "目標ID（detail/update/delete/progress 使用）"},
                "steps": {"type": "string", "description": "執行步驟或新狀態值（add 時可選填，update_status 填新狀態如 in_progress/completed，set_progress 填0-100）"},
                "priority": {"type": "string", "enum": ["low","normal","high","critical"], "description": "優先級（預設normal）"}
            },
            "required": ["action"]
        }
    },
    {
        "name": "auto_trade",
        "description": "加密貨幣交易。查詢Binance即時價格、帳戶餘額、買賣下單、查看掛單。需要設定 BINANCE_KEY 和 BINANCE_SECRET 才能下單。",
        "input_schema": {
            "type": "object",
            "properties": {
                "action": {"type": "string", "enum": ["price","balance","buy","sell","open_orders"],
                           "description": "price=查價格, balance=查餘額, buy=買入, sell=賣出, open_orders=查掛單"},
                "symbol": {"type": "string", "description": "交易對（如 BTCUSDT, ETHUSDT, BNBUSDT）"},
                "amount": {"type": "number", "description": "交易數量（buy/sell 使用）"},
                "price": {"type": "number", "description": "限價（limit order 使用，market order 不需要）"},
                "order_type": {"type": "string", "enum": ["market","limit"], "description": "訂單類型（預設market市價）"}
            },
            "required": ["action"]
        }
    },
    {
        "name": "knowledge_base",
        "description": "知識庫管理。儲存重要知識/資料/筆記、全文搜尋、查看詳情、統計。當用戶要儲存資訊、查找之前存的知識時使用。",
        "input_schema": {
            "type": "object",
            "properties": {
                "action": {"type": "string", "enum": ["add","search","get","list","delete","stats"],
                           "description": "add=新增知識, search=搜尋, get=查看詳情, list=列出所有, delete=刪除, stats=統計"},
                "content": {"type": "string", "description": "知識內容（add 使用）"},
                "query": {"type": "string", "description": "搜尋關鍵字（search 使用）"},
                "tag": {"type": "string", "description": "標籤分類（add 使用，如 科技/財經/個人）"},
                "kb_id": {"type": "string", "description": "知識ID（get/delete 使用）"}
            },
            "required": ["action"]
        }
    },
    {
        "name": "emotion_detect",
        "description": "情緒偵測分析。從文字分析情緒狀態、從螢幕/圖片分析臉部情緒。當用戶想了解情緒狀態、分析對話情緒時使用。",
        "input_schema": {
            "type": "object",
            "properties": {
                "action": {"type": "string", "enum": ["from_text","from_face"],
                           "description": "from_text=文字情緒分析, from_face=臉部情緒分析（截圖或指定圖片）"},
                "text": {"type": "string", "description": "要分析的文字（from_text 使用）"},
                "image_path": {"type": "string", "description": "圖片路徑（from_face 使用，不填則截圖）"}
            },
            "required": ["action"]
        }
    },
    {
        "name": "voice_id",
        "description": "聲紋辨識。登記聲紋特徵、辨識說話者身份、管理已登記的聲紋資料庫。需要麥克風或音訊檔案。",
        "input_schema": {
            "type": "object",
            "properties": {
                "action": {"type": "string", "enum": ["enroll","identify","list","delete"],
                           "description": "enroll=登記新聲紋, identify=辨識說話者, list=列出已登記, delete=刪除聲紋"},
                "name": {"type": "string", "description": "人物名稱（enroll/delete 使用）"},
                "audio_path": {"type": "string", "description": "音訊檔路徑（identify 使用）"},
                "duration": {"type": "number", "description": "錄音秒數（enroll 使用，預設5秒）"}
            },
            "required": ["action"]
        }
    },
    {
        "name": "pentest",
        "description": "安全評估工具（僅用於自己的網路）。埠掃描、SSL憑證檢查、HTTP安全標頭分析、網頁漏洞掃描、密碼強度審計。",
        "input_schema": {
            "type": "object",
            "properties": {
                "action": {"type": "string", "enum": ["port_scan","ssl_check","http_headers","vuln_scan","password_audit"],
                           "description": "port_scan=掃描開放埠, ssl_check=SSL憑證, http_headers=安全標頭, vuln_scan=網頁漏洞, password_audit=密碼強度"},
                "target": {"type": "string", "description": "目標IP/域名/URL（password_audit 填密碼，多個用逗號分隔）"},
                "port_range": {"type": "string", "description": "埠範圍（port_scan 使用，如 1-1000 或 80,443,3389）"},
                "timeout": {"type": "number", "description": "連線逾時秒數（預設2）"}
            },
            "required": ["action","target"]
        }
    },
    {
        "name": "proactive_alert",
        "description": "主動預警系統。設定條件（CPU/記憶體/加密貨幣價格/新聞關鍵字）超過閾值時自動通知。啟動後背景持續監控。",
        "input_schema": {
            "type": "object",
            "properties": {
                "action": {"type": "string", "enum": ["add","list","delete","toggle","start_all"],
                           "description": "add=新增預警, list=查看所有, delete=刪除, toggle=啟用/停用, start_all=啟動所有監控"},
                "name": {"type": "string", "description": "預警名稱"},
                "condition": {"type": "string", "enum": ["cpu_above","memory_above","price_above","price_below","news_keyword"],
                              "description": "條件類型"},
                "threshold": {"type": "string", "description": "閾值（cpu/memory填百分比數字，price填價格）"},
                "target": {"type": "string", "description": "監控目標（price用交易對如BTCUSDT，news_keyword用關鍵字）"},
                "interval": {"type": "number", "description": "檢查間隔秒數（預設60）"}
            },
            "required": ["action"]
        }
    },
    {
        "name": "multi_deploy",
        "description": "多機器部署。將bot部署到遠端伺服器、查看遠端bot狀態、同步技能更新到遠端機器、查看遠端日誌。需要SSH存取權限。",
        "input_schema": {
            "type": "object",
            "properties": {
                "action": {"type": "string", "enum": ["deploy","status","sync","log"],
                           "description": "deploy=完整部署, status=查看狀態, sync=同步技能並重啟, log=查看日誌"},
                "remote_host": {"type": "string", "description": "遠端主機IP或域名"},
                "remote_user": {"type": "string", "description": "SSH用戶名"},
                "remote_pass": {"type": "string", "description": "SSH密碼"},
                "remote_path": {"type": "string", "description": "遠端部署路徑（預設 /tmp/niu_bot）"}
            },
            "required": ["action","remote_host","remote_user","remote_pass"]
        }
    },
    {
        "name": "self_benchmark",
        "description": "自我評估。檢查所有功能健康狀態、統計已部署技能數量、查看記憶體統計。當用戶問bot目前狀態、健康檢查時使用。",
        "input_schema": {
            "type": "object",
            "properties": {
                "action": {"type": "string", "enum": ["run","skill_count","memory_stats"],
                           "description": "run=完整健康報告, skill_count=技能清單, memory_stats=記憶統計"}
            },
            "required": ["action"]
        },
        "cache_control": {"type": "ephemeral"}
    },
    {
        "name": "think_as",
        "description": "用指定人物的思維框架分析問題。載入蒸餾好的心智模型、決策啟發式、表達DNA，用那個人的角度思考。當用戶說「用XX的角度分析」「XX會怎麼看」「用XX的思維」時使用。可用人物：elon-musk(馬斯克)、warren-buffett(巴菲特)、jensen-huang(黃仁勳)、morris-chang(張忠謀)。",
        "input_schema": {
            "type": "object",
            "properties": {
                "person": {"type": "string", "description": "人物名稱，如 elon-musk, warren-buffett, jensen-huang, morris-chang"},
                "question": {"type": "string", "description": "要用這個人的角度分析的問題"},
                "list_available": {"type": "boolean", "description": "設為true列出所有可用人物"}
            },
            "required": ["person", "question"]
        }
    }
]


def execute_think_as(person: str, question: str, list_available: bool = False) -> str:
    """載入蒸餾好的人物思維框架，用該人物角度分析問題"""
    from pathlib import Path
    skills_dir = Path(__file__).parent / "skills"

    if list_available:
        available = [f.stem for f in skills_dir.glob("*.md") if f.stem != "colleague-niuma"]
        return "🧠 可用的思維框架：\n" + "\n".join(f"- {name}" for name in available) if available else "目前沒有可用的思維框架"

    # 找對應的 skill 檔案
    slug = person.lower().strip().replace(" ", "-").replace("_", "-")
    # 別名對照
    aliases = {
        "馬斯克": "elon-musk", "musk": "elon-musk", "elon": "elon-musk",
        "巴菲特": "warren-buffett", "buffett": "warren-buffett", "warren": "warren-buffett",
        "黃仁勳": "jensen-huang", "jensen": "jensen-huang", "huang": "jensen-huang",
        "張忠謀": "morris-chang", "morris": "morris-chang", "chang": "morris-chang",
    }
    slug = aliases.get(slug, slug)
    skill_path = skills_dir / f"{slug}.md"

    if not skill_path.exists():
        available = [f.stem for f in skills_dir.glob("*.md") if f.stem != "colleague-niuma"]
        return f"找不到「{person}」的思維框架。可用的有：{', '.join(available)}"

    # 讀取 skill 內容
    skill_content = skill_path.read_text(encoding="utf-8")

    # 用 Claude 載入框架分析問題
    resp = client.messages.create(
        model="claude-sonnet-4-6",
        max_tokens=1500,
        system=f"你現在要完全用以下人物的思維框架來分析問題。不是模仿說話方式，是用他的心智模型和決策規則來推導。\n\n{skill_content}",
        messages=[{"role": "user", "content": question}]
    )
    return f"🧠 {person} 的視角：\n\n{resp.content[0].text}"


def calc_rsi(closes, period=14):
    """計算 RSI 指標"""
    deltas = closes.diff().dropna()
    gains = deltas.clip(lower=0)
    losses = -deltas.clip(upper=0)
    avg_gain = gains.rolling(period).mean().iloc[-1]
    avg_loss = losses.rolling(period).mean().iloc[-1]
    if avg_loss == 0:
        return 100.0
    rs = avg_gain / avg_loss
    return round(100 - (100 / (1 + rs)), 1)


def fetch_stock(symbol: str, period: str = "1mo") -> str:
    try:
        import yfinance as yf
        ticker = yf.Ticker(symbol)
        info = ticker.info
        # 取較長歷史以便計算 RSI
        hist = ticker.history(period="3mo")

        if hist.empty:
            return f"找不到「{symbol}」的股票數據，請確認代號是否正確。"

        name = info.get("longName") or info.get("shortName") or symbol
        currency = info.get("currency", "")
        current = hist["Close"].iloc[-1]
        prev = hist["Close"].iloc[-2] if len(hist) > 1 else current
        change = current - prev
        change_pct = (change / prev) * 100 if prev else 0
        arrow = "▲" if change >= 0 else "▼"
        volume = hist["Volume"].iloc[-1]
        avg_volume = hist["Volume"].tail(20).mean()
        volume_ratio = volume / avg_volume if avg_volume else 1

        # 技術指標
        ma5 = hist["Close"].tail(5).mean()
        ma20 = hist["Close"].tail(20).mean()
        ma60 = hist["Close"].tail(60).mean() if len(hist) >= 60 else hist["Close"].mean()
        rsi = calc_rsi(hist["Close"]) if len(hist) >= 15 else None

        # 趨勢
        if ma5 > ma20 > ma60:
            trend = "強勢多頭（MA5>MA20>MA60）📈"
        elif ma5 < ma20 < ma60:
            trend = "強勢空頭（MA5<MA20<MA60）📉"
        elif ma5 > ma20:
            trend = "短線偏多（MA5>MA20）🔼"
        else:
            trend = "短線偏空（MA5<MA20）🔽"

        # RSI 解讀
        rsi_note = ""
        if rsi is not None:
            if rsi >= 80:
                rsi_note = "（嚴重超買 ⚠️）"
            elif rsi >= 70:
                rsi_note = "（超買區間）"
            elif rsi <= 20:
                rsi_note = "（嚴重超賣 💡）"
            elif rsi <= 30:
                rsi_note = "（超賣區間）"
            else:
                rsi_note = "（中性）"

        # 基本面
        market_cap = info.get("marketCap")
        pe_ratio = info.get("trailingPE")
        week52_high = info.get("fiftyTwoWeekHigh")
        week52_low = info.get("fiftyTwoWeekLow")
        high_period = hist["High"].tail(20).max()
        low_period = hist["Low"].tail(20).min()

        # 距離高低點百分比
        pct_from_high = ((current - high_period) / high_period * 100) if high_period else 0
        pct_from_low = ((current - low_period) / low_period * 100) if low_period else 0

        result = (
            f"📊 {name} ({symbol})\n"
            f"💰 現價：{current:.2f} {currency}  {arrow} {abs(change):.2f} ({change_pct:+.2f}%)\n"
            f"📦 成交量：{volume:,}（均量 {volume_ratio:.1f}x）\n"
            f"\n── 技術指標 ──\n"
            f"MA5：{ma5:.2f}　MA20：{ma20:.2f}　MA60：{ma60:.2f}\n"
            f"趨勢：{trend}\n"
        )

        if rsi is not None:
            result += f"RSI(14)：{rsi}{rsi_note}\n"

        result += (
            f"近20日高點：{high_period:.2f}（距高 {pct_from_high:.1f}%）\n"
            f"近20日低點：{low_period:.2f}（距低 +{pct_from_low:.1f}%）\n"
        )

        if week52_high and week52_low:
            result += f"52週高低：{week52_low:.2f} ~ {week52_high:.2f}\n"

        result += "\n── 基本面 ──\n"
        if market_cap:
            mc_str = f"{market_cap/1e12:.2f}兆" if market_cap >= 1e12 else f"{market_cap/1e8:.0f}億"
            result += f"市值：{mc_str} {currency}\n"
        if pe_ratio:
            result += f"本益比：{pe_ratio:.1f}\n"

        return result.strip()

    except Exception as e:
        return f"查詢「{symbol}」失敗：{str(e)}"


def fetch_ashare(code: str, period: str = "1mo") -> str:
    """A股（滬深）/ 港股查詢，自動判斷市場並加後綴"""
    try:
        import yfinance as yf
        code = code.strip().lstrip("0" * 0)  # 保留原始代號

        # 判斷市場：6位數字 → A股（6開頭=上海.SS，其他=深圳.SZ）；4位以下 → 港股.HK
        if code.isdigit():
            if len(code) == 6:
                symbol = f"{code}.SS" if code.startswith("6") else f"{code}.SZ"
                market = "A股"
            else:
                symbol = f"{code.zfill(4)}.HK"
                market = "港股"
        else:
            symbol = code  # 使用者直接提供後綴
            market = "未知"

        ticker = yf.Ticker(symbol)
        info = ticker.info
        hist = ticker.history(period="3mo")

        if hist.empty:
            return f"找不到「{code}」的數據，請確認代號是否正確。"

        name = info.get("longName") or info.get("shortName") or code
        currency = info.get("currency", "")
        current = hist["Close"].iloc[-1]
        prev = hist["Close"].iloc[-2] if len(hist) > 1 else current
        change = current - prev
        change_pct = (change / prev * 100) if prev else 0
        arrow = "▲" if change >= 0 else "▼"
        volume = hist["Volume"].iloc[-1]

        ma5  = hist["Close"].tail(5).mean()
        ma20 = hist["Close"].tail(20).mean()
        ma60 = hist["Close"].tail(60).mean() if len(hist) >= 60 else hist["Close"].mean()
        rsi  = calc_rsi(hist["Close"]) if len(hist) >= 15 else None

        if ma5 > ma20 > ma60:
            trend = "強勢多頭 📈"
        elif ma5 < ma20 < ma60:
            trend = "強勢空頭 📉"
        elif ma5 > ma20:
            trend = "短線偏多 🔼"
        else:
            trend = "短線偏空 🔽"

        rsi_note = ""
        if rsi is not None:
            if rsi >= 80:   rsi_note = "（嚴重超買 ⚠️）"
            elif rsi >= 70: rsi_note = "（超買）"
            elif rsi <= 20: rsi_note = "（嚴重超賣 💡）"
            elif rsi <= 30: rsi_note = "（超賣）"
            else:           rsi_note = "（中性）"

        week52_high = info.get("fiftyTwoWeekHigh")
        week52_low  = info.get("fiftyTwoWeekLow")
        market_cap  = info.get("marketCap")
        pe_ratio    = info.get("trailingPE")

        result = (
            f"🇨🇳 {name}（{market} {code}）\n"
            f"💰 現價：{current:.2f} {currency}  {arrow} {abs(change):.2f} ({change_pct:+.2f}%)\n"
            f"📦 成交量：{volume:,}\n"
            f"\n── 技術指標 ──\n"
            f"MA5：{ma5:.2f}　MA20：{ma20:.2f}　MA60：{ma60:.2f}\n"
            f"趨勢：{trend}\n"
        )
        if rsi is not None:
            result += f"RSI(14)：{rsi}{rsi_note}\n"
        if week52_high and week52_low:
            result += f"52週高低：{week52_low:.2f} ~ {week52_high:.2f}\n"
        result += "\n── 基本面 ──\n"
        if market_cap:
            mc_str = f"{market_cap/1e12:.2f}兆" if market_cap >= 1e12 else f"{market_cap/1e8:.0f}億"
            result += f"市值：{mc_str} {currency}\n"
        if pe_ratio:
            result += f"本益比：{pe_ratio:.1f}\n"

        return result.strip()
    except Exception as e:
        return f"查詢「{code}」失敗：{e}"


def fetch_cn_news(source: str = "all", count: int = 5) -> str:
    """抓取中國大陸新聞 RSS"""
    try:
        import feedparser
        count = min(max(count, 1), 10)
        feeds = {
            "xinhua":  ("新華社",   "https://feeds.xinhuanet.com/news/world/rss"),
            "people":  ("人民網",   "http://www.people.com.cn/rss/politics.xml"),
            "36kr":    ("36氪",     "https://36kr.com/feed"),
            "caixin":  ("財新網",   "https://www.caixin.com/rss/home.xml"),
        }
        # 備用可靠來源
        fallback_feeds = {
            "xinhua":  ("新華社",   "https://rsshub.app/xinhua/world"),
            "people":  ("人民網",   "https://rsshub.app/people/politics"),
            "36kr":    ("36氪",     "https://rsshub.app/36kr/news/latest"),
            "caixin":  ("財新網",   "https://rsshub.app/caixin/blog"),
        }
        sources = list(feeds.keys()) if source == "all" else [source]
        results = []
        for src in sources:
            if src not in feeds:
                continue
            label, url = feeds[src]
            try:
                feed = feedparser.parse(url)
                items = feed.entries[:count]
                if not items:
                    # 嘗試備用
                    label2, url2 = fallback_feeds.get(src, (label, url))
                    feed = feedparser.parse(url2)
                    items = feed.entries[:count]
                if not items:
                    results.append(f"📰 {label}：暫無資料")
                    continue
                lines = [f"📰 {label}"]
                for i, entry in enumerate(items, 1):
                    title = entry.get("title", "無標題")
                    lines.append(f"{i}. {title}")
                results.append("\n".join(lines))
            except Exception:
                results.append(f"📰 {label}：抓取失敗")
        return "\n\n".join(results) if results else "無法取得中國新聞"
    except Exception as e:
        return f"中國新聞查詢失敗：{e}"


def fetch_china_search(query: str, category: str = "其他", count: int = 6) -> str:
    """全方位中國大陸資訊搜尋：旅遊/美食/文化/戲劇/演員/工作等"""
    try:
        count = min(max(count, 1), 10)
        results = []

        # 1. Google News（中文簡體，抓最新新聞/資訊）
        try:
            import feedparser
            news_query = query
            # 依分類補充關鍵字讓搜尋更精準
            category_hints = {
                "旅遊": f"{query} 旅遊攻略 景點",
                "美食": f"{query} 美食 餐廳 推薦",
                "文化風俗": f"{query} 文化 習俗 傳統",
                "戲劇影視": f"{query} 電視劇 電影 劇情",
                "演員明星": f"{query} 演員 明星 近況",
                "工作生活": f"{query} 工作 薪資 生活",
                "城市介紹": f"{query} 城市 介紹 特色",
                "歷史": f"{query} 歷史 背景",
                "科技": f"{query} 科技 技術",
                "新聞時事": f"{query} 最新 新聞",
            }
            news_query = category_hints.get(category, query)
            url = f"https://news.google.com/rss/search?q={urllib.parse.quote(news_query)}&hl=zh-Hans&gl=CN&ceid=CN:zh-Hans"
            feed = feedparser.parse(url)
            if feed.entries:
                lines = [f"📰 Google 新聞（{category}）"]
                for i, entry in enumerate(feed.entries[:min(count, 5)], 1):
                    title = entry.get("title", "").split(" - ")[0]  # 去掉媒體名稱
                    pub = entry.get("published", "")[:16]
                    lines.append(f"{i}. {title}（{pub}）")
                results.append("\n".join(lines))
        except Exception:
            pass

        # 2. DuckDuckGo 搜尋（zh-cn 地區，含陸網資料）
        try:
            from ddgs import DDGS
            ddg_results = []
            with DDGS() as ddgs:
                for r in ddgs.text(query, region="zh-cn", max_results=count):
                    title = r.get("title", "")
                    body  = r.get("body", "")[:150]
                    href  = r.get("href", "")
                    ddg_results.append(f"• {title}\n  {body}\n  {href}")
            if ddg_results:
                results.append(f"🔍 網路搜尋結果\n" + "\n\n".join(ddg_results))
        except Exception:
            pass

        # 3. Wikipedia 中文（適合文化/歷史/人物類）
        if category in ("文化風俗", "歷史", "演員明星", "城市介紹", "戲劇影視"):
            try:
                wiki_url = f"https://zh.wikipedia.org/api/rest_v1/page/summary/{urllib.parse.quote(query.split()[0])}"
                resp = requests.get(wiki_url, timeout=8)
                if resp.status_code == 200:
                    data = resp.json()
                    extract = data.get("extract", "")
                    title_w = data.get("title", "")
                    if extract:
                        results.append(f"📖 Wikipedia：{title_w}\n{extract[:500]}")
            except Exception:
                pass

        if not results:
            return f"找不到「{query}」的相關資訊，請嘗試更換關鍵字"

        return "\n\n" + "─" * 25 + "\n\n".join(results)

    except Exception as e:
        return f"中國資訊搜尋失敗：{e}"


def fetch_institutional(symbol: str = "", date: str = "") -> str:
    """台股三大法人買賣超"""
    try:
        import datetime
        if not date:
            date = datetime.date.today().strftime("%Y%m%d")
        headers = {"User-Agent": "Mozilla/5.0"}

        if symbol:
            # 個股三大法人
            url = f"https://www.twse.com.tw/rwd/zh/fund/T86?response=json&date={date}&selectType=ALL"
            resp = requests.get(url, timeout=10, headers=headers)
            data = resp.json()
            if data.get("stat") != "OK":
                return f"查無資料（{date}，可能為非交易日）"
            rows = data.get("data", [])
            target = None
            for row in rows:
                if str(row[0]).strip() == str(symbol).strip():
                    target = row
                    break
            if not target:
                return f"找不到 {symbol} 的三大法人資料"
            foreign = int(target[4].replace(",", "").replace("+", ""))
            trust   = int(target[10].replace(",", "").replace("+", ""))
            dealer  = int(target[13].replace(",", "").replace("+", "")) if len(target) > 13 else 0
            total   = foreign + trust + dealer
            arrow = lambda v: "▲" if v >= 0 else "▼"
            return (
                f"📊 {symbol} 三大法人（{date[:4]}/{date[4:6]}/{date[6:]}）\n"
                f"外資：{arrow(foreign)} {abs(foreign):,} 張\n"
                f"投信：{arrow(trust)} {abs(trust):,} 張\n"
                f"自營：{arrow(dealer)} {abs(dealer):,} 張\n"
                f"合計：{arrow(total)} {abs(total):,} 張"
            )
        else:
            # 市場整體
            url = f"https://www.twse.com.tw/rwd/zh/fund/BFI82U?response=json&type=day&dayDate={date}"
            resp = requests.get(url, timeout=10, headers=headers)
            data = resp.json()
            if data.get("stat") != "OK":
                return f"查無三大法人整體資料（{date}，可能為非交易日）"
            rows = data.get("data", [])
            lines = [f"📊 台股三大法人整體買賣超（{date[:4]}/{date[4:6]}/{date[6:]}）\n"]
            for row in rows:
                name = row[0]
                buy  = row[1].replace(",", "")
                sell = row[2].replace(",", "")
                diff = row[3].replace(",", "")
                val  = int(diff.replace("+", "")) if diff.replace("+","").replace("-","").isdigit() else 0
                arrow = "▲" if val >= 0 else "▼"
                lines.append(f"{name}：{arrow} {diff} 元")
            return "\n".join(lines)
    except Exception as e:
        return f"三大法人查詢失敗：{e}"


def fetch_sector(market: str = "us") -> str:
    """產業類股表現"""
    try:
        import yfinance as yf
        if market == "us":
            sectors = {
                "科技": "XLK", "金融": "XLF", "醫療": "XLV", "能源": "XLE",
                "工業": "XLI", "消費必需": "XLP", "消費選擇": "XLY",
                "公用事業": "XLU", "材料": "XLB", "通訊": "XLC", "房地產": "XLRE"
            }
        else:
            sectors = {
                "半導體": "00891.TW", "金融": "0055.TW", "航運": "00895.TW",
                "電動車": "00893.TW", "ESG": "00878.TW", "高息": "00919.TW",
                "科技": "0052.TW", "傳產": "0054.TW"
            }

        results = []
        for name, sym in sectors.items():
            try:
                hist = yf.Ticker(sym).history(period="2d")
                if len(hist) >= 2:
                    chg = (hist["Close"].iloc[-1] / hist["Close"].iloc[-2] - 1) * 100
                    arrow = "▲" if chg >= 0 else "▼"
                    results.append((chg, f"{arrow} {name}：{chg:+.2f}%"))
            except Exception:
                pass

        results.sort(key=lambda x: x[0], reverse=True)
        market_label = "美股" if market == "us" else "台股"
        lines = [f"🏭 {market_label}產業類股今日表現\n"]
        for _, line in results:
            lines.append(line)
        if results:
            lines.append(f"\n最強：{results[0][1].split('：')[0].strip()}")
            lines.append(f"最弱：{results[-1][1].split('：')[0].strip()}")
        return "\n".join(lines)
    except Exception as e:
        return f"類股查詢失敗：{e}"


def fetch_commodity(items: list = None) -> str:
    """黃金/原油/原物料報價"""
    try:
        import yfinance as yf
        commodity_map = {
            "gold":   ("黃金",   "GC=F",  "USD/盎司"),
            "oil":    ("WTI原油", "CL=F",  "USD/桶"),
            "silver": ("白銀",   "SI=F",  "USD/盎司"),
            "copper": ("銅",     "HG=F",  "USD/磅"),
            "natgas": ("天然氣", "NG=F",  "USD/MMBtu"),
            "wheat":  ("小麥",   "ZW=F",  "USd/英斗"),
            "corn":   ("玉米",   "ZC=F",  "USd/英斗"),
        }
        if not items or "all" in items:
            items = list(commodity_map.keys())

        lines = ["🛢 大宗商品報價\n"]
        for key in items:
            if key not in commodity_map:
                continue
            name, sym, unit = commodity_map[key]
            try:
                hist = yf.Ticker(sym).history(period="2d")
                if hist.empty:
                    continue
                price = hist["Close"].iloc[-1]
                if len(hist) >= 2:
                    chg = (price / hist["Close"].iloc[-2] - 1) * 100
                    arrow = "▲" if chg >= 0 else "▼"
                    lines.append(f"{name}：{price:.2f} {unit}  {arrow} {chg:+.2f}%")
                else:
                    lines.append(f"{name}：{price:.2f} {unit}")
            except Exception:
                pass
        return "\n".join(lines)
    except Exception as e:
        return f"商品報價失敗：{e}"


def fetch_bond_yield() -> str:
    """美國公債殖利率"""
    try:
        import yfinance as yf
        bonds = {
            "3個月": "^IRX",
            "2年":   "2YY=F",
            "5年":   "^FVX",
            "10年":  "^TNX",
            "30年":  "^TYX",
        }
        lines = ["🏦 美國公債殖利率\n"]
        yields = {}
        for label, sym in bonds.items():
            try:
                hist = yf.Ticker(sym).history(period="2d")
                if hist.empty:
                    continue
                val = hist["Close"].iloc[-1]
                yields[label] = val
                if len(hist) >= 2:
                    chg = val - hist["Close"].iloc[-2]
                    arrow = "▲" if chg >= 0 else "▼"
                    lines.append(f"{label}：{val:.3f}%  {arrow} {chg:+.3f}%")
                else:
                    lines.append(f"{label}：{val:.3f}%")
            except Exception:
                pass

        # 利差分析
        if "2年" in yields and "10年" in yields:
            spread = yields["10年"] - yields["2年"]
            curve = "正斜率（景氣正常）" if spread > 0 else "倒掛（衰退警訊 ⚠️）"
            lines.append(f"\n10Y-2Y 利差：{spread:+.3f}%  {curve}")
        if "3個月" in yields and "10年" in yields:
            spread2 = yields["10年"] - yields["3個月"]
            curve2 = "正常" if spread2 > 0 else "倒掛 ⚠️"
            lines.append(f"10Y-3M 利差：{spread2:+.3f}%  {curve2}")

        return "\n".join(lines)
    except Exception as e:
        return f"殖利率查詢失敗：{e}"


def fetch_dividend_calendar(symbol: str) -> str:
    """除權息資訊"""
    try:
        import yfinance as yf
        ticker = yf.Ticker(symbol)
        info = ticker.info
        name = info.get("longName") or info.get("shortName") or symbol
        currency = info.get("currency", "")

        div_yield    = info.get("dividendYield", 0) or 0
        div_rate     = info.get("dividendRate", 0) or 0
        ex_date      = info.get("exDividendDate")
        last_div     = info.get("lastDividendValue", 0) or 0
        payout_ratio = info.get("payoutRatio", 0) or 0
        price        = info.get("currentPrice") or info.get("regularMarketPrice", 0)

        lines = [f"💰 {name} ({symbol}) 除權息資訊\n"]
        if div_rate:
            lines.append(f"年配息：{div_rate:.4f} {currency}")
        if div_yield:
            lines.append(f"殖利率：{div_yield*100:.2f}%")
        if last_div:
            lines.append(f"上次配息：{last_div:.4f} {currency}")
        if ex_date:
            import datetime
            ex_dt = datetime.datetime.fromtimestamp(ex_date).strftime("%Y-%m-%d")
            lines.append(f"除息日：{ex_dt}")
        if payout_ratio:
            lines.append(f"配息率：{payout_ratio*100:.1f}%")

        # 歷史配息紀錄
        try:
            divs = ticker.dividends
            if divs is not None and not divs.empty:
                lines.append("\n近5次配息：")
                for dt, val in divs.tail(5).iloc[::-1].items():
                    lines.append(f"  {str(dt)[:10]}：{val:.4f} {currency}")
        except Exception:
            pass

        return "\n".join(lines) if len(lines) > 1 else f"{symbol} 無配息資料"
    except Exception as e:
        return f"除權息查詢失敗：{e}"


def fetch_stock_screener(criteria: str, market: str = "us") -> str:
    """選股篩選器（用 Claude 解讀條件 + yfinance 驗證）"""
    try:
        import yfinance as yf

        # 美股用 S&P500 成分股，台股用常見大型股
        if market == "us":
            # 取 S&P500 部分成分股做示範
            candidates = [
                "AAPL","MSFT","GOOGL","AMZN","NVDA","META","TSLA","BRK-B","JPM","V",
                "XOM","UNH","JNJ","WMT","MA","PG","HD","CVX","MRK","ABBV",
                "KO","PEP","BAC","PFE","AVGO","COST","TMO","MCD","ABT","CRM",
                "ACN","LIN","DHR","TXN","NEE","QCOM","PM","HON","IBM","GE",
                "ORCL","AMGN","SBUX","CAT","INTU","AMD","ISRG","NOW","MDLZ","AXP"
            ]
        else:
            candidates = [
                "2330.TW","2317.TW","2454.TW","2412.TW","2308.TW","2303.TW",
                "2881.TW","2882.TW","2886.TW","2891.TW","2002.TW","1301.TW",
                "0050.TW","0056.TW","00878.TW","00919.TW","2603.TW","2609.TW",
                "3711.TW","2379.TW","3008.TW","2395.TW","4938.TW","2376.TW"
            ]

        results = []
        # 解析條件關鍵字
        want_high_div   = any(k in criteria for k in ["殖利率", "配息", "dividend"])
        want_low_pe     = any(k in criteria for k in ["本益比", "PE", "pe"])
        want_high_roe   = any(k in criteria for k in ["ROE", "roe", "股東權益"])
        want_top_gain   = any(k in criteria for k in ["漲最多", "漲幅", "漲停", "上漲"])
        want_large_cap  = any(k in criteria for k in ["市值", "大型", "large"])

        for sym in candidates[:30]:  # 限制查詢數量避免超時
            try:
                info = yf.Ticker(sym).info
                name = info.get("shortName", sym)
                price = info.get("currentPrice") or info.get("regularMarketPrice", 0)
                pe    = info.get("trailingPE", 0) or 0
                div_y = (info.get("dividendYield", 0) or 0) * 100
                roe   = (info.get("returnOnEquity", 0) or 0) * 100
                cap   = info.get("marketCap", 0) or 0
                chg   = info.get("regularMarketChangePercent", 0) or 0

                score = 0
                if want_high_div and div_y > 3:   score += div_y
                if want_low_pe   and 0 < pe < 20: score += (20 - pe)
                if want_high_roe and roe > 15:     score += roe / 10
                if want_top_gain:                  score += chg
                if want_large_cap and cap > 1e11:  score += 1

                if score > 0:
                    results.append((score, sym, name, price, pe, div_y, roe, chg, cap))
            except Exception:
                pass

        results.sort(reverse=True)
        lines = [f"🔎 選股結果（{criteria}）\n"]
        for i, (_, sym, name, price, pe, div_y, roe, chg, cap) in enumerate(results[:10], 1):
            cap_str = f"{cap/1e12:.1f}兆" if cap >= 1e12 else f"{cap/1e8:.0f}億"
            line = f"{i}. {name}（{sym}）  {chg:+.1f}%"
            if div_y: line += f"  殖利率{div_y:.1f}%"
            if pe:    line += f"  PE{pe:.0f}"
            if roe:   line += f"  ROE{roe:.0f}%"
            lines.append(line)
        if not results:
            lines.append("找不到符合條件的股票，試著調整條件")
        return "\n".join(lines)
    except Exception as e:
        return f"選股篩選失敗：{e}"


def fetch_margin_trading(symbol: str, date: str = "") -> str:
    """台股融資融券餘額"""
    try:
        import datetime
        if not date:
            date = datetime.date.today().strftime("%Y%m%d")
        # TWSE 融資融券 API
        url = f"https://www.twse.com.tw/rwd/zh/marginTrading/MI_MARGN?response=json&date={date}&selectType=ALL"
        resp = requests.get(url, timeout=10, headers={"User-Agent": "Mozilla/5.0"})
        data = resp.json()
        if data.get("stat") != "OK":
            return f"查無融資融券資料（{date}，可能為非交易日）"
        rows = data.get("data", [])
        target = None
        for row in rows:
            if str(row[0]).strip() == str(symbol).strip():
                target = row
                break
        if not target:
            return f"找不到 {symbol} 的融資融券資料"

        # 欄位：代號, 名稱, 融資買進, 融資賣出, 融資現金償還, 融資餘額, 融資限額,
        #        融券賣出, 融券買進, 融券現券償還, 融券餘額, 融券限額, 資券互抵
        name     = target[1]
        loan_bal = target[5].replace(",", "")   # 融資餘額（千股）
        short_bal= target[10].replace(",", "")  # 融券餘額（千股）
        loan_buy = target[2].replace(",", "")
        loan_sell= target[3].replace(",", "")
        short_sell=target[7].replace(",", "")
        short_buy= target[8].replace(",", "")

        return (
            f"📋 {name}（{symbol}）融資融券（{date[:4]}/{date[4:6]}/{date[6:]}）\n\n"
            f"── 融資（散戶多單）──\n"
            f"餘額：{loan_bal} 千股\n"
            f"今買：{loan_buy} 千股　今賣：{loan_sell} 千股\n\n"
            f"── 融券（放空）──\n"
            f"餘額：{short_bal} 千股\n"
            f"今賣：{short_sell} 千股　今買：{short_buy} 千股\n\n"
            f"資券比：{int(loan_bal)/(int(short_bal) if int(short_bal) > 0 else 1):.1f}x"
            if loan_bal.isdigit() and short_bal.isdigit() else ""
        )
    except Exception as e:
        return f"融資融券查詢失敗：{e}"


def fetch_options(symbol: str, expiry: str = "") -> str:
    """選擇權鏈資料"""
    try:
        import yfinance as yf
        ticker = yf.Ticker(symbol)
        expirations = ticker.options
        if not expirations:
            return f"{symbol} 無選擇權資料"

        exp = expiry if expiry in expirations else expirations[0]
        chain = ticker.option_chain(exp)
        calls = chain.calls
        puts  = chain.puts

        price = ticker.info.get("currentPrice") or ticker.info.get("regularMarketPrice", 0)

        # 取最接近現價的 5 個履約價
        calls = calls.iloc[(calls["strike"] - price).abs().argsort()[:5]].sort_values("strike")
        puts  = puts.iloc[(puts["strike"]  - price).abs().argsort()[:5]].sort_values("strike")

        lines = [f"📈 {symbol} 選擇權（到期：{exp}，現價：{price:.2f}）\n"]
        lines.append("── Call（買權）──")
        for _, row in calls.iterrows():
            iv = f"IV {row.get('impliedVolatility', 0)*100:.0f}%" if row.get("impliedVolatility") else ""
            oi = f"OI {row.get('openInterest', 0):,}" if row.get("openInterest") else ""
            lines.append(f"  履約 {row['strike']:.0f}：{row.get('lastPrice', 0):.2f}  {iv}  {oi}")

        lines.append("\n── Put（賣權）──")
        for _, row in puts.iterrows():
            iv = f"IV {row.get('impliedVolatility', 0)*100:.0f}%" if row.get("impliedVolatility") else ""
            oi = f"OI {row.get('openInterest', 0):,}" if row.get("openInterest") else ""
            lines.append(f"  履約 {row['strike']:.0f}：{row.get('lastPrice', 0):.2f}  {iv}  {oi}")

        return "\n".join(lines)
    except Exception as e:
        return f"選擇權查詢失敗：{e}"


def fetch_futures(items: list = None) -> str:
    """主要期貨報價"""
    try:
        import yfinance as yf
        futures_map = {
            "sp500":  ("S&P500期貨",    "ES=F"),
            "nasdaq": ("那斯達克期貨",  "NQ=F"),
            "dow":    ("道瓊期貨",      "YM=F"),
            "gold":   ("黃金期貨",      "GC=F"),
            "oil":    ("WTI原油期貨",   "CL=F"),
            "taiex":  ("台指期",        "TAIEX.TW"),
        }
        if not items or "all" in items:
            items = list(futures_map.keys())

        lines = ["📊 期貨報價\n"]
        for key in items:
            if key not in futures_map:
                continue
            name, sym = futures_map[key]
            try:
                hist = yf.Ticker(sym).history(period="2d")
                if hist.empty:
                    continue
                price = hist["Close"].iloc[-1]
                if len(hist) >= 2:
                    chg = (price / hist["Close"].iloc[-2] - 1) * 100
                    arrow = "▲" if chg >= 0 else "▼"
                    lines.append(f"{name}：{price:,.2f}  {arrow} {chg:+.2f}%")
                else:
                    lines.append(f"{name}：{price:,.2f}")
            except Exception:
                pass
        return "\n".join(lines)
    except Exception as e:
        return f"期貨查詢失敗：{e}"


def fetch_ipo(count: int = 10) -> str:
    """近期 IPO 行事曆"""
    try:
        import feedparser
        count = min(count, 20)
        results = []

        # 用 Google 新聞搜尋 IPO 資訊
        url = f"https://news.google.com/rss/search?q=IPO+新股+上市&hl=zh-Hant&gl=TW&ceid=TW:zh-Hant"
        feed = feedparser.parse(url)
        lines = ["🆕 近期 IPO / 新股資訊\n"]
        for entry in feed.entries[:count]:
            title = entry.get("title", "").split(" - ")[0]
            pub   = entry.get("published", "")[:16]
            lines.append(f"• {title}（{pub}）")

        # 補充美股 IPO（用 DuckDuckGo）
        try:
            from ddgs import DDGS
            with DDGS() as ddgs:
                for r in ddgs.text("upcoming IPO 2026 stock market", region="us-en", max_results=5):
                    title = r.get("title", "")
                    body  = r.get("body", "")[:100]
                    lines.append(f"🇺🇸 {title}\n   {body}")
        except Exception:
            pass

        return "\n".join(lines) if len(lines) > 1 else "暫無 IPO 資訊"
    except Exception as e:
        return f"IPO 查詢失敗：{e}"


def fetch_backtest(symbol: str, strategy: str = "ma_cross", period: str = "2y") -> str:
    """回測投資策略"""
    try:
        import yfinance as yf
        import numpy as np

        ticker = yf.Ticker(symbol)
        hist = ticker.history(period=period)
        if hist.empty or len(hist) < 30:
            return f"找不到 {symbol} 的歷史資料"

        name = ticker.info.get("shortName") or symbol
        close = hist["Close"]
        n = len(close)

        if strategy == "buy_hold":
            total_ret = (close.iloc[-1] / close.iloc[0] - 1) * 100
            annual_ret = ((close.iloc[-1] / close.iloc[0]) ** (252 / n) - 1) * 100
            max_dd = ((close / close.cummax()) - 1).min() * 100
            return (
                f"📈 {name}（{symbol}）買進持有回測（{period}）\n\n"
                f"起始價：{close.iloc[0]:.2f}\n"
                f"結束價：{close.iloc[-1]:.2f}\n"
                f"總報酬：{total_ret:+.2f}%\n"
                f"年化報酬：{annual_ret:+.2f}%\n"
                f"最大回撤：{max_dd:.2f}%"
            )

        elif strategy == "ma_cross":
            ma5  = close.rolling(5).mean()
            ma20 = close.rolling(20).mean()
            signal = (ma5 > ma20).astype(int)
            signal_prev = signal.shift(1)
            buy_signals  = (signal == 1) & (signal_prev == 0)
            sell_signals = (signal == 0) & (signal_prev == 1)

            trades = []
            buy_price = None
            for i in range(len(close)):
                if buy_signals.iloc[i] and buy_price is None:
                    buy_price = close.iloc[i]
                elif sell_signals.iloc[i] and buy_price is not None:
                    ret = (close.iloc[i] / buy_price - 1) * 100
                    trades.append(ret)
                    buy_price = None

            if not trades:
                return f"{symbol} 在 {period} 內無均線交叉訊號"

            wins = sum(1 for t in trades if t > 0)
            total_ret = sum(trades)
            avg_ret   = total_ret / len(trades)
            win_rate  = wins / len(trades) * 100
            max_dd    = ((close / close.cummax()) - 1).min() * 100

            return (
                f"📊 {name}（{symbol}）MA5穿MA20 策略回測（{period}）\n\n"
                f"交易次數：{len(trades)} 次\n"
                f"勝率：{win_rate:.1f}%\n"
                f"平均每筆報酬：{avg_ret:+.2f}%\n"
                f"累計報酬：{total_ret:+.2f}%\n"
                f"最大回撤：{max_dd:.2f}%\n\n"
                f"買進持有同期：{(close.iloc[-1]/close.iloc[0]-1)*100:+.2f}%"
            )

        elif strategy == "dca":
            # 每月第一個交易日買進固定金額
            monthly = close.resample("MS").first()
            shares = 0
            cost   = 0
            invest_per_month = 10000  # 每月投入 10000 元
            for price in monthly:
                shares += invest_per_month / price
                cost   += invest_per_month
            final_value = shares * close.iloc[-1]
            total_ret   = (final_value / cost - 1) * 100

            return (
                f"📅 {name}（{symbol}）定期定額回測（{period}，每月10000元）\n\n"
                f"總投入：{cost:,.0f} 元\n"
                f"最終市值：{final_value:,.0f} 元\n"
                f"總報酬：{total_ret:+.2f}%\n"
                f"累積股數：{shares:.2f} 股\n"
                f"平均成本：{cost/shares:.2f} 元"
            )

        return "不支援的策略"
    except Exception as e:
        return f"回測失敗：{e}"


def fetch_global_market() -> str:
    try:
        import yfinance as yf
        markets = {
            "🇺🇸 S&P500": "^GSPC", "🇺🇸 那斯達克": "^IXIC", "🇺🇸 道瓊": "^DJI",
            "🇹🇼 台股": "^TWII", "🇭🇰 恆生": "^HSI", "🇯🇵 日經": "^N225",
            "🇰🇷 韓股": "^KS11", "🇬🇧 英國": "^FTSE", "🇩🇪 德國": "^GDAXI",
            "🇫🇷 法國": "^FCHI", "🌏 上證": "000001.SS",
        }
        lines = ["🌍 全球市場概覽\n"]
        for name, sym in markets.items():
            try:
                hist = yf.Ticker(sym).history(period="2d")
                if len(hist) >= 2:
                    chg = (hist["Close"].iloc[-1] / hist["Close"].iloc[-2] - 1) * 100
                    price = hist["Close"].iloc[-1]
                    arrow = "▲" if chg >= 0 else "▼"
                    lines.append(f"{name}：{price:,.2f}  {arrow} {chg:+.2f}%")
            except Exception:
                pass
        return "\n".join(lines)
    except Exception as e:
        return f"全球市場查詢失敗：{e}"


def fetch_economic_calendar(count: int = 10) -> str:
    try:
        import feedparser
        results = []
        # Investing.com RSS 與 Google 新聞
        urls = [
            ("https://news.google.com/rss/search?q=CPI+非農+Fed利率+GDP+經濟數據&hl=zh-Hant&gl=TW&ceid=TW:zh-Hant", "財經日曆"),
            ("https://news.google.com/rss/search?q=economic+calendar+CPI+nonfarm+Fed+GDP&hl=en-US&gl=US&ceid=US:en", "Economic Calendar"),
        ]
        lines = ["📅 重要經濟數據行事曆\n"]
        seen = set()
        for url, label in urls:
            feed = feedparser.parse(url)
            for entry in feed.entries:
                title = entry.get("title", "").split(" - ")[0]
                pub = entry.get("published", "")[:16]
                key = title[:30]
                if key not in seen:
                    seen.add(key)
                    lines.append(f"• {title}（{pub}）")
                if len(lines) > count + 1:
                    break
            if len(lines) > count + 1:
                break
        lines.append("\n💡 重點關注：CPI（通膨）、非農就業（NFP）、Fed利率決議、GDP、PPI")
        return "\n".join(lines)
    except Exception as e:
        return f"經濟日曆查詢失敗：{e}"


def fetch_earnings_calendar(days: int = 7) -> str:
    try:
        import feedparser, datetime
        lines = [f"📊 未來 {days} 天財報行事曆\n"]
        url = "https://news.google.com/rss/search?q=earnings+report+quarterly+results&hl=en-US&gl=US&ceid=US:en"
        feed = feedparser.parse(url)
        count = 0
        for entry in feed.entries:
            title = entry.get("title", "").split(" - ")[0]
            pub = entry.get("published", "")[:16]
            lines.append(f"• {title}（{pub}）")
            count += 1
            if count >= 10:
                break

        # 補充重點大型股財報
        import yfinance as yf
        watchlist = ["AAPL", "MSFT", "GOOGL", "AMZN", "NVDA", "META", "TSLA"]
        lines.append("\n重點股財報日：")
        for sym in watchlist:
            try:
                cal = yf.Ticker(sym).calendar
                if cal and "Earnings Date" in cal:
                    ed = cal["Earnings Date"]
                    if hasattr(ed, '__iter__'):
                        ed = list(ed)[0]
                    lines.append(f"  {sym}：{str(ed)[:10]}")
            except Exception:
                pass
        return "\n".join(lines)
    except Exception as e:
        return f"財報日曆查詢失敗：{e}"


def fetch_analyst_ratings(symbol: str) -> str:
    try:
        import yfinance as yf
        ticker = yf.Ticker(symbol)
        name = ticker.info.get("shortName") or symbol

        upgrades = ticker.upgrades_downgrades
        if upgrades is None or upgrades.empty:
            return f"{symbol} 無分析師評級資料"

        upgrades = upgrades.reset_index()
        recent = upgrades.head(10)
        lines = [f"📋 {name}（{symbol}）分析師評級\n"]
        for _, row in recent.iterrows():
            date = str(row.get("GradeDate", ""))[:10]
            firm = row.get("Firm", "")
            action = row.get("Action", "")
            to_grade = row.get("ToGrade", "")
            from_grade = row.get("FromGrade", "")
            emoji = "📈" if "up" in str(action).lower() or "initiat" in str(action).lower() else ("📉" if "down" in str(action).lower() else "➡️")
            lines.append(f"{emoji} {date} {firm}：{from_grade}→{to_grade}（{action}）")

        price = ticker.info.get("currentPrice", 0)
        target = ticker.info.get("targetMeanPrice", 0)
        if price and target:
            upside = (target / price - 1) * 100
            lines.append(f"\n分析師均價目標：{target:.2f}（現價 {price:.2f}，空間 {upside:+.1f}%）")
        return "\n".join(lines)
    except Exception as e:
        return f"分析師評級查詢失敗：{e}"


def fetch_short_interest(symbol: str) -> str:
    try:
        import yfinance as yf
        info = yf.Ticker(symbol).info
        name = info.get("shortName") or symbol

        short_pct = info.get("shortPercentOfFloat", 0) or 0
        short_ratio = info.get("shortRatio", 0) or 0
        shares_short = info.get("sharesShort", 0) or 0
        shares_out = info.get("sharesOutstanding", 1) or 1

        if short_pct >= 0.20:
            risk = "極高空頭壓力（軋空機會大 ⚠️）"
        elif short_pct >= 0.10:
            risk = "高空頭比率（需留意）"
        elif short_pct >= 0.05:
            risk = "中等空頭比率"
        else:
            risk = "低空頭比率（市場偏多）"

        return (
            f"🩳 {name}（{symbol}）空頭資料\n\n"
            f"做空比率：{short_pct*100:.2f}%\n"
            f"回補天數（Short Ratio）：{short_ratio:.1f} 天\n"
            f"借券賣出股數：{shares_short:,}\n"
            f"風險評估：{risk}"
        )
    except Exception as e:
        return f"空頭資料查詢失敗：{e}"


def fetch_correlation(symbols: list, period: str = "1y") -> str:
    try:
        import yfinance as yf
        import pandas as pd
        data = {}
        for sym in symbols[:5]:
            hist = yf.Ticker(sym).history(period=period)
            if not hist.empty:
                data[sym] = hist["Close"].pct_change().dropna()

        if len(data) < 2:
            return "至少需要 2 個有效的股票代號"

        df = pd.DataFrame(data).dropna()
        corr = df.corr()

        lines = [f"📊 相關性矩陣（{period}）\n"]
        syms = list(corr.columns)
        header = "      " + "  ".join(f"{s:>8}" for s in syms)
        lines.append(header)
        for s1 in syms:
            row = f"{s1:>6}"
            for s2 in syms:
                val = corr.loc[s1, s2]
                row += f"  {val:>8.3f}"
            lines.append(row)

        lines.append("\n💡 相關係數：1.0=完全正相關，0=無關，-1.0=完全負相關")
        lines.append("分散風險建議選相關係數 < 0.5 的資產")

        # 找最低相關配對
        pairs = []
        for i, s1 in enumerate(syms):
            for s2 in syms[i+1:]:
                pairs.append((corr.loc[s1, s2], s1, s2))
        pairs.sort()
        if pairs:
            v, s1, s2 = pairs[0]
            lines.append(f"最低相關：{s1} & {s2}（{v:.3f}）← 最佳分散組合")
        return "\n".join(lines)
    except Exception as e:
        return f"相關性分析失敗：{e}"


def fetch_risk_metrics(symbol: str, period: str = "1y") -> str:
    try:
        import yfinance as yf
        import numpy as np

        ticker = yf.Ticker(symbol)
        hist = ticker.history(period=period)
        bench = yf.Ticker("^GSPC").history(period=period)
        name = ticker.info.get("shortName") or symbol

        if hist.empty or bench.empty:
            return f"找不到 {symbol} 資料"

        ret = hist["Close"].pct_change().dropna()
        bench_ret = bench["Close"].pct_change().dropna()
        aligned = ret.align(bench_ret, join="inner")
        ret, bench_ret = aligned[0], aligned[1]

        # Beta
        cov = np.cov(ret, bench_ret)
        beta = cov[0][1] / cov[1][1] if cov[1][1] != 0 else 0

        # 年化報酬 & 波動率
        annual_ret = ret.mean() * 252 * 100
        annual_vol = ret.std() * (252 ** 0.5) * 100

        # 夏普比率（無風險利率 5%）
        rf = 0.05
        sharpe = (annual_ret/100 - rf) / (annual_vol/100) if annual_vol != 0 else 0

        # 最大回撤
        cumret = (1 + ret).cumprod()
        max_dd = ((cumret / cumret.cummax()) - 1).min() * 100

        # VaR 95%
        var_95 = np.percentile(ret, 5) * 100

        return (
            f"⚖️ {name}（{symbol}）風險指標（{period}）\n\n"
            f"Beta（市場敏感度）：{beta:.2f}{'（高波動）' if beta>1.2 else '（低波動）' if beta<0.8 else '（接近大盤）'}\n"
            f"年化報酬：{annual_ret:+.2f}%\n"
            f"年化波動率：{annual_vol:.2f}%\n"
            f"夏普比率：{sharpe:.2f}{'（優秀）' if sharpe>1 else '（尚可）' if sharpe>0.5 else '（偏低）'}\n"
            f"最大回撤：{max_dd:.2f}%\n"
            f"VaR 95%（單日）：{var_95:.2f}%"
        )
    except Exception as e:
        return f"風險指標計算失敗：{e}"


def fetch_money_flow(symbol: str) -> str:
    try:
        import yfinance as yf
        ticker = yf.Ticker(symbol)
        hist = ticker.history(period="20d")
        info = ticker.info
        name = info.get("shortName") or symbol

        if hist.empty:
            return f"找不到 {symbol} 資料"

        # 資金流向：(收盤-開盤)/當日振幅 * 成交量 估算買賣壓
        typical = (hist["High"] + hist["Low"] + hist["Close"]) / 3
        raw_mf = typical * hist["Volume"]
        pos_mf = raw_mf[hist["Close"] >= hist["Open"]].tail(10).sum()
        neg_mf = raw_mf[hist["Close"] < hist["Open"]].tail(10).sum()
        mfi = 100 - (100 / (1 + pos_mf / neg_mf)) if neg_mf else 100

        today_vol = hist["Volume"].iloc[-1]
        avg_vol = hist["Volume"].tail(20).mean()
        vol_ratio = today_vol / avg_vol if avg_vol else 1

        price = hist["Close"].iloc[-1]
        chg = (price / hist["Close"].iloc[-2] - 1) * 100 if len(hist) > 1 else 0

        flow = "淨流入 📥" if mfi > 55 else "淨流出 📤" if mfi < 45 else "中性"
        return (
            f"💰 {name}（{symbol}）資金流向\n\n"
            f"今日漲跌：{chg:+.2f}%\n"
            f"成交量：{today_vol:,}（均量 {vol_ratio:.1f}x）\n"
            f"資金流向指標（MFI）：{mfi:.1f} → {flow}\n"
            f"近10日正向資金：{pos_mf/1e8:.1f}億\n"
            f"近10日負向資金：{neg_mf/1e8:.1f}億"
        )
    except Exception as e:
        return f"資金流向查詢失敗：{e}"


def fetch_concept_stocks(theme: str) -> str:
    try:
        # 台股概念股資料庫
        concepts = {
            "AI": ["2330.TW","2303.TW","2454.TW","3711.TW","2379.TW","2308.TW","6488.TW","3017.TW","2382.TW","5274.TW"],
            "電動車": ["2308.TW","1590.TW","2207.TW","6239.TW","1802.TW","2049.TW","3665.TW","1537.TW","6227.TW","2371.TW"],
            "軍工": ["1323.TW","2409.TW","2348.TW","8112.TW","1536.TW","2634.TW","6245.TW","1513.TW"],
            "低軌衛星": ["3045.TW","2230.TW","6438.TW","3413.TW","2365.TW","3508.TW","6285.TW","4306.TW"],
            "半導體": ["2330.TW","2303.TW","2454.TW","2308.TW","3711.TW","5347.TW","6770.TW","3034.TW","2379.TW","4919.TW"],
            "5G": ["2412.TW","3045.TW","2498.TW","2356.TW","3231.TW","6488.TW","2439.TW","3293.TW"],
            "儲能": ["1907.TW","1504.TW","6409.TW","3481.TW","2023.TW","6121.TW","1590.TW"],
            "DRAM": ["2303.TW","3450.TW","4967.TW","3260.TW","2408.TW"],
            "CoWoS": ["2330.TW","6235.TW","3711.TW","2454.TW","3036.TW","8046.TW","6415.TW"],
            "矽光子": ["2330.TW","3008.TW","2454.TW","6510.TW","3081.TW"],
            "機器人": ["2308.TW","1590.TW","2049.TW","1537.TW","3665.TW","6288.TW","2382.TW"],
            "航運": ["2603.TW","2609.TW","2615.TW","2610.TW","2618.TW","5608.TW"],
            "金融": ["2881.TW","2882.TW","2886.TW","2891.TW","2884.TW","2892.TW","5876.TW"],
        }

        # 模糊比對
        import yfinance as yf
        matched_key = None
        for key in concepts:
            if key in theme or theme in key:
                matched_key = key
                break

        if not matched_key:
            # 用 DuckDuckGo 搜尋
            try:
                from ddgs import DDGS
                with DDGS() as ddgs:
                    results = list(ddgs.text(f"台股 {theme} 概念股 相關股票", region="zh-tw", max_results=5))
                lines = [f"🏭 {theme} 概念股\n（搜尋結果）"]
                for r in results:
                    lines.append(f"• {r.get('title','')}\n  {r.get('body','')[:100]}")
                return "\n\n".join(lines)
            except Exception:
                return f"找不到「{theme}」概念股，支援：{'、'.join(concepts.keys())}"

        syms = concepts[matched_key]
        lines = [f"🏭 {matched_key} 概念股（台股）\n"]
        for sym in syms:
            try:
                hist = yf.Ticker(sym).history(period="2d")
                info = yf.Ticker(sym).info
                name = info.get("shortName") or sym.replace(".TW","")
                if len(hist) >= 2:
                    chg = (hist["Close"].iloc[-1] / hist["Close"].iloc[-2] - 1) * 100
                    price = hist["Close"].iloc[-1]
                    arrow = "▲" if chg >= 0 else "▼"
                    lines.append(f"{sym.replace('.TW','')} {name}：{price:.1f}  {arrow}{chg:+.1f}%")
            except Exception:
                lines.append(sym.replace(".TW",""))
        return "\n".join(lines)
    except Exception as e:
        return f"概念股查詢失敗：{e}"


def fetch_crypto_depth(coin: str = "bitcoin") -> str:
    try:
        coin_map = {
            "btc": "bitcoin", "eth": "ethereum", "sol": "solana",
            "bnb": "binancecoin", "xrp": "ripple", "doge": "dogecoin",
        }
        coin_id = coin_map.get(coin.lower(), coin.lower())

        # CoinGecko 詳細資料
        resp = requests.get(
            f"https://api.coingecko.com/api/v3/coins/{coin_id}"
            "?localization=false&tickers=false&community_data=true&developer_data=false",
            timeout=10
        )
        data = resp.json()
        if "error" in data:
            return f"找不到幣種「{coin}」"

        md = data["market_data"]
        name = data["name"]
        sym = data["symbol"].upper()
        price = md["current_price"]["usd"]
        ch24 = md.get("price_change_percentage_24h") or 0
        ch7d = md.get("price_change_percentage_7d") or 0
        mcap = md["market_cap"]["usd"]
        vol = md["total_volume"]["usd"]
        dom = data.get("market_cap_percentage", {}).get(coin_id.split("-")[0], 0)

        # 資金費率（用 DuckDuckGo 補充）
        funding_note = ""
        try:
            from ddgs import DDGS
            with DDGS() as ddgs:
                r = list(ddgs.text(f"{sym} funding rate perpetual", max_results=1))
                if r:
                    funding_note = f"\n資金費率參考：{r[0].get('body','')[:80]}"
        except Exception:
            pass

        return (
            f"🔗 {name}（{sym}）鏈上深度\n\n"
            f"現價：${price:,.4f}\n"
            f"24h：{ch24:+.2f}%　7d：{ch7d:+.2f}%\n"
            f"市值：${mcap/1e9:.2f}B\n"
            f"24h 成交量：${vol/1e9:.2f}B\n"
            f"社群分數：{data.get('community_score',0):.0f}/100\n"
            f"開發者分數：{data.get('developer_score',0):.0f}/100"
            + funding_note
        )
    except Exception as e:
        return f"加密幣深度查詢失敗：{e}"


def fetch_drip_calculator(symbol: str, shares: float, years: int = 10, monthly_invest: float = 0) -> str:
    try:
        import yfinance as yf
        ticker = yf.Ticker(symbol)
        info = ticker.info
        name = info.get("shortName") or symbol
        price = info.get("currentPrice") or info.get("regularMarketPrice") or 0
        div_yield = (info.get("dividendYield") or 0)
        div_rate = info.get("dividendRate") or (price * div_yield)
        payout_freq = 4  # 假設季配

        if not price:
            return f"找不到 {symbol} 的股價資料"

        total_shares = shares
        total_invested = shares * price
        annual_divs = []

        for year in range(1, years + 1):
            # 年度股息 → 再買股
            annual_div = total_shares * div_rate
            new_shares = annual_div / price if price > 0 else 0
            total_shares += new_shares

            # 每月定期追加
            if monthly_invest > 0:
                monthly_shares = (monthly_invest * 12) / price
                total_shares += monthly_shares
                total_invested += monthly_invest * 12

            annual_divs.append(annual_div)

        final_value = total_shares * price
        total_div = sum(annual_divs)
        total_return = (final_value / total_invested - 1) * 100 if total_invested > 0 else 0

        lines = [
            f"💰 {name}（{symbol}）DRIP 股息再投資試算\n",
            f"初始：{shares:.0f} 股 × {price:.2f} = {shares*price:,.0f} 元",
        ]
        if monthly_invest > 0:
            lines.append(f"每月追加：{monthly_invest:,.0f} 元")
        lines += [
            f"殖利率：{div_yield*100:.2f}%　每股年配：{div_rate:.4f}",
            f"\n{years} 年後：",
            f"持股數：{total_shares:,.1f} 股",
            f"總市值：{final_value:,.0f} 元",
            f"總投入：{total_invested:,.0f} 元",
            f"累積股息：{total_div:,.0f} 元",
            f"總報酬：{total_return:+.1f}%",
            f"年配息（第{years}年）：{annual_divs[-1]:,.0f} 元",
        ]
        return "\n".join(lines)
    except Exception as e:
        return f"DRIP 試算失敗：{e}"


def fetch_forex_chart(pair: str, period: str = "3mo") -> str:
    try:
        import yfinance as yf
        ticker = yf.Ticker(pair)
        hist = ticker.history(period="6mo")
        info = ticker.info

        if hist.empty:
            return f"找不到匯率「{pair}」，請確認格式（如 USDTWD=X）"

        name = info.get("shortName") or pair
        close = hist["Close"]
        current = close.iloc[-1]
        prev = close.iloc[-2] if len(close) > 1 else current
        chg = (current / prev - 1) * 100

        ma5  = close.tail(5).mean()
        ma20 = close.tail(20).mean()
        ma60 = close.tail(60).mean() if len(close) >= 60 else close.mean()
        rsi  = calc_rsi(close) if len(close) >= 15 else None

        if ma5 > ma20 > ma60:   trend = "強勢升值 📈"
        elif ma5 < ma20 < ma60: trend = "強勢貶值 📉"
        elif ma5 > ma20:        trend = "短線偏強 🔼"
        else:                   trend = "短線偏弱 🔽"

        from ta.volatility import BollingerBands
        bb = BollingerBands(close)
        bb_upper = bb.bollinger_hband().iloc[-1]
        bb_lower = bb.bollinger_lband().iloc[-1]
        bb_pos = "近上軌（偏強）" if current >= bb_upper * 0.99 else "近下軌（偏弱）" if current <= bb_lower * 1.01 else "通道中間"

        result = (
            f"💱 {name} 技術分析\n"
            f"現值：{current:.4f}  {'▲' if chg >= 0 else '▼'} {chg:+.2f}%\n\n"
            f"MA5：{ma5:.4f}　MA20：{ma20:.4f}　MA60：{ma60:.4f}\n"
            f"趨勢：{trend}\n"
            f"布林：{bb_pos}\n"
        )
        if rsi:
            rsi_note = "超買" if rsi >= 70 else "超賣" if rsi <= 30 else "中性"
            result += f"RSI(14)：{rsi}（{rsi_note}）\n"
        result += (
            f"\n52週高：{hist['High'].max():.4f}\n"
            f"52週低：{hist['Low'].min():.4f}"
        )
        return result
    except Exception as e:
        return f"外匯技術分析失敗：{e}"


def fetch_warrant(underlying: str) -> str:
    try:
        # TWSE 權證查詢
        url = f"https://www.twse.com.tw/rwd/zh/warrant/MNO01?response=json&stockNo={underlying}"
        resp = requests.get(url, timeout=10, headers={"User-Agent": "Mozilla/5.0"})
        data = resp.json()

        if data.get("stat") != "OK" or not data.get("data"):
            # 改用搜尋提示
            try:
                from ddgs import DDGS
                with DDGS() as ddgs:
                    r = list(ddgs.text(f"台股 {underlying} 權證 認購 認售", region="zh-tw", max_results=3))
                lines = [f"🎫 {underlying} 相關權證資訊\n"]
                for item in r:
                    lines.append(f"• {item.get('title','')}\n  {item.get('body','')[:100]}")
                return "\n\n".join(lines)
            except Exception:
                return f"查無 {underlying} 的權證資料，請至台灣證券交易所查詢"

        rows = data.get("data", [])[:10]
        lines = [f"🎫 {underlying} 相關權證（前10筆）\n"]
        for row in rows:
            w_code = row[0]
            w_name = row[1]
            w_type = "認購" if "購" in w_name else "認售"
            strike = row[3] if len(row) > 3 else "-"
            exp = row[4] if len(row) > 4 else "-"
            lines.append(f"{w_code} {w_name}（{w_type}）　履約{strike}　到期{exp}")
        return "\n".join(lines)
    except Exception as e:
        return f"權證查詢失敗：{e}"


def fetch_portfolio_risk(holdings: list, period: str = "1y") -> str:
    try:
        import yfinance as yf
        import numpy as np
        import pandas as pd

        if not holdings or not isinstance(holdings, list):
            return "⚠️ 請提供持股清單，格式：[{\"symbol\": \"AAPL\", \"weight\": 50}, ...]"
        if isinstance(holdings[0], str):
            holdings = [{"symbol": s, "weight": 1} for s in holdings]
        weights = [h.get("weight", 1) if isinstance(h, dict) else 1 for h in holdings]
        symbols = [h.get("symbol", h) if isinstance(h, dict) else str(h) for h in holdings]

        # 標準化權重
        total_w = sum(weights)
        if total_w == 0:
            return "⚠️ 所有持股權重為 0，請設定有效權重"
        weights = [w / total_w for w in weights]

        data = {}
        for sym in symbols:
            hist = yf.Ticker(sym).history(period=period)
            if not hist.empty:
                data[sym] = hist["Close"].pct_change().dropna()

        if not data:
            return "無法取得任何股票資料"

        df = pd.DataFrame(data).dropna()
        valid_syms = list(df.columns)
        valid_weights = [weights[symbols.index(s)] for s in valid_syms]
        valid_weights = [w / sum(valid_weights) for w in valid_weights]

        # 組合報酬
        portfolio_ret = (df * valid_weights).sum(axis=1)
        annual_ret = portfolio_ret.mean() * 252 * 100
        annual_vol = portfolio_ret.std() * (252 ** 0.5) * 100
        sharpe = (annual_ret/100 - 0.05) / (annual_vol/100) if annual_vol else 0
        max_dd = ((1 + portfolio_ret).cumprod() / (1 + portfolio_ret).cumprod().cummax() - 1).min() * 100
        var_95 = np.percentile(portfolio_ret, 5) * 100

        # 相關性
        corr = df.corr()

        lines = [f"📊 投資組合風險分析（{period}）\n"]
        lines.append("持倉配置：")
        for sym, w in zip(valid_syms, valid_weights):
            lines.append(f"  {sym}：{w*100:.1f}%")

        lines += [
            f"\n── 組合風險指標 ──",
            f"年化報酬：{annual_ret:+.2f}%",
            f"年化波動率：{annual_vol:.2f}%",
            f"夏普比率：{sharpe:.2f}",
            f"最大回撤：{max_dd:.2f}%",
            f"VaR 95%（單日）：{var_95:.2f}%",
        ]

        if len(valid_syms) >= 2:
            lines.append("\n── 相關性（越低越分散）──")
            for i, s1 in enumerate(valid_syms):
                for s2 in valid_syms[i+1:]:
                    v = corr.loc[s1, s2]
                    note = "⚠️高度相關" if v > 0.7 else "✅低相關" if v < 0.3 else ""
                    lines.append(f"  {s1} & {s2}：{v:.3f} {note}")

        return "\n".join(lines)
    except Exception as e:
        return f"投資組合風險分析失敗：{e}"


def fetch_retirement_calculator(current_age: int, current_savings: float, monthly_save: float,
                                 retire_age: int = 65, annual_return: float = 6.0,
                                 monthly_expense: float = 50000) -> str:
    try:
        years = retire_age - current_age
        if years <= 0:
            return "退休年齡必須大於目前年齡"
        r_monthly = annual_return / 100 / 12
        # 現有資產複利成長
        future_current = current_savings * 10000 * ((1 + r_monthly) ** (years * 12))
        # 每月儲蓄複利成長（年金終值）
        if r_monthly > 0:
            future_monthly = monthly_save * (((1 + r_monthly) ** (years * 12) - 1) / r_monthly)
        else:
            future_monthly = monthly_save * years * 12
        total_at_retire = future_current + future_monthly
        # 退休後可用年數（假設活到85歲）
        retire_years = 85 - retire_age
        total_needed = monthly_expense * 12 * retire_years
        gap = total_at_retire - total_needed
        status = "✅ 達標" if gap >= 0 else "⚠️ 不足"

        lines = [
            f"🏖️ 退休規劃試算\n",
            f"目前年齡：{current_age} 歲　預計退休：{retire_age} 歲",
            f"距退休：{years} 年　預期報酬：{annual_return}%/年",
            f"",
            f"── 退休時預估資產 ──",
            f"現有資產成長至：{future_current/10000:.0f} 萬元",
            f"累積儲蓄成長至：{future_monthly/10000:.0f} 萬元",
            f"退休時總資產：{total_at_retire/10000:.0f} 萬元",
            f"",
            f"── 退休所需評估（活至85歲）──",
            f"月生活費：{monthly_expense:,.0f} 元",
            f"退休後需要：{total_needed/10000:.0f} 萬元",
            f"缺口/剩餘：{gap/10000:+.0f} 萬元　{status}",
        ]
        if gap < 0:
            extra = abs(gap) / (years * 12)
            lines.append(f"\n每月需額外存：{extra:,.0f} 元 才能達標")
        return "\n".join(lines)
    except Exception as e:
        return f"退休試算失敗：{e}"


def fetch_loan_calculator(principal: float, annual_rate: float, years: int,
                           loan_type: str = "等額本息") -> str:
    try:
        p = principal * 10000
        r = annual_rate / 100 / 12
        n = years * 12
        if n <= 0:
            return "⚠️ 貸款年數必須大於 0"
        total_interest = 0
        lines = [f"🏦 貸款試算（{loan_type}）\n",
                 f"貸款金額：{principal:.0f} 萬元",
                 f"年利率：{annual_rate}%　期數：{n} 期（{years} 年）\n"]
        if loan_type == "等額本息":
            if r > 0:
                payment = p * r * (1 + r) ** n / ((1 + r) ** n - 1)
            else:
                payment = p / n
            total_pay = payment * n
            total_interest = total_pay - p
            lines += [
                f"── 等額本息 ──",
                f"每月還款：{payment:,.0f} 元",
                f"總還款額：{total_pay/10000:.2f} 萬元",
                f"總利息：{total_interest/10000:.2f} 萬元",
            ]
        else:  # 等額本金
            principal_payment = p / n
            first_payment = principal_payment + p * r
            last_payment = principal_payment + principal_payment * r
            total_interest = sum(principal_payment * r * (n - i) for i in range(n))
            lines += [
                f"── 等額本金 ──",
                f"每期本金：{principal_payment:,.0f} 元",
                f"第1期還款：{first_payment:,.0f} 元",
                f"最後1期還款：{last_payment:,.0f} 元",
                f"總利息：{total_interest/10000:.2f} 萬元",
                f"總還款額：{(p + total_interest)/10000:.2f} 萬元",
            ]
        return "\n".join(lines)
    except Exception as e:
        return f"貸款試算失敗：{e}"


def fetch_compound_calculator(principal: float, annual_rate: float, years: int,
                               monthly_add: float = 0, compound_freq: int = 12) -> str:
    try:
        compound_freq = int(compound_freq) if not isinstance(compound_freq, int) else compound_freq
        if compound_freq <= 0:
            compound_freq = 12
        r = annual_rate / 100 / compound_freq
        n = compound_freq * years
        # 本金複利
        future_principal = principal * (1 + r) ** n
        # 每月定期投入（按月計算）
        if monthly_add > 0:
            r_m = annual_rate / 100 / 12
            n_m = years * 12
            future_monthly = monthly_add * (((1 + r_m) ** n_m - 1) / r_m) if r_m > 0 else monthly_add * n_m
        else:
            future_monthly = 0
        total = future_principal + future_monthly
        total_invest = principal + monthly_add * years * 12
        profit = total - total_invest
        lines = [
            f"📈 複利計算器\n",
            f"本金：{principal:,.0f} 元",
            f"年化報酬：{annual_rate}%　期間：{years} 年",
            f"每月追加：{monthly_add:,.0f} 元",
            f"",
            f"── 試算結果 ──",
            f"本金複利成長：{future_principal:,.0f} 元",
        ]
        if monthly_add > 0:
            lines.append(f"追加投入成長：{future_monthly:,.0f} 元")
        lines += [
            f"期末總資產：{total:,.0f} 元",
            f"總投入成本：{total_invest:,.0f} 元",
            f"獲利：{profit:,.0f} 元（{profit/total_invest*100:.1f}%）",
            f"",
            f"── 72法則 ──",
            f"資產翻倍需：{72/annual_rate:.1f} 年（年報酬 {annual_rate}%）",
        ]
        return "\n".join(lines)
    except Exception as e:
        return f"複利計算失敗：{e}"


def fetch_asset_allocation(age: int, risk_level: str, goal: str = "退休",
                            investment_horizon: int = None) -> str:
    try:
        horizon = investment_horizon or max(65 - age, 5)
        # 基本股債比（110法則）
        base_stock = min(110 - age, 90)
        if risk_level == "保守":
            stock = max(base_stock - 20, 10)
        elif risk_level == "積極":
            stock = min(base_stock + 15, 90)
        else:
            stock = base_stock
        bond = max(100 - stock - 10, 0)
        cash = 100 - stock - bond

        # 細分配置
        tw_stock = round(stock * 0.4)
        us_stock = round(stock * 0.35)
        intl_stock = stock - tw_stock - us_stock
        tw_bond = round(bond * 0.3)
        us_bond = bond - tw_bond

        lines = [
            f"📊 資產配置建議\n",
            f"年齡：{age} 歲　風險偏好：{risk_level}",
            f"目標：{goal}　投資期間：{horizon} 年\n",
            f"── 大類配置 ──",
            f"股票：{stock}%　債券：{bond}%　現金：{cash}%\n",
            f"── 細分建議 ──",
            f"台股（0050/0056）：{tw_stock}%",
            f"美股（VTI/VOO）：{us_stock}%",
            f"國際股（VEA/VWO）：{intl_stock}%",
            f"台灣公債/ETF：{tw_bond}%",
            f"美債（BND/TLT）：{us_bond}%",
            f"現金/定存：{cash}%\n",
            f"── 再平衡建議 ──",
            f"每年或偏離5%以上時再平衡",
            f"隨年齡增長逐步降低股票比例",
        ]
        if risk_level == "保守":
            lines.append("\n⚠️ 保守型：優先保本，適合距退休較近者")
        elif risk_level == "積極":
            lines.append("\n💡 積極型：承受較大波動換取長期成長，需有10年以上視野")
        return "\n".join(lines)
    except Exception as e:
        return f"資產配置建議失敗：{e}"


def fetch_tw_tax_calculator(dividend_income: float, other_income: float = 0,
                             tax_bracket: float = None, sell_amount: float = 0) -> str:
    try:
        # 健保補充保費（2.11%，超過2萬才扣）
        nhi_surcharge = dividend_income * 0.0211 if dividend_income >= 20000 else 0
        # 股利可抵減稅額（8.5%，上限8萬）
        tax_credit = min(dividend_income * 0.085, 80000)
        # 分離課稅（28%）
        separate_tax = dividend_income * 0.28
        # 合併申報
        if tax_bracket:
            total_income = dividend_income + other_income
            combined_tax = total_income * (tax_bracket / 100) - tax_credit
            combined_tax = max(combined_tax, 0)
        else:
            combined_tax = None
        # 證交稅（0.3%）
        securities_tax = sell_amount * 0.003

        lines = [
            f"💰 台股稅務試算\n",
            f"股利所得：{dividend_income:,.0f} 元",
        ]
        if other_income:
            lines.append(f"其他收入：{other_income:,.0f} 元")
        lines += [
            f"",
            f"── 健保補充保費（2.11%）──",
            f"補充保費：{nhi_surcharge:,.0f} 元",
            f"",
            f"── 方案一：分離課稅（28%）──",
            f"應繳稅額：{separate_tax:,.0f} 元",
            f"稅後股利：{dividend_income - separate_tax - nhi_surcharge:,.0f} 元",
        ]
        if combined_tax is not None:
            lines += [
                f"",
                f"── 方案二：合併申報（稅率{tax_bracket}%）──",
                f"可抵減稅額：{tax_credit:,.0f} 元",
                f"應繳稅額：{combined_tax:,.0f} 元",
                f"稅後股利：{dividend_income - combined_tax - nhi_surcharge:,.0f} 元",
                f"",
                f"建議：{'合併申報' if combined_tax < separate_tax else '分離課稅'} 節稅 {abs(separate_tax - combined_tax):,.0f} 元",
            ]
        if sell_amount > 0:
            lines += [
                f"",
                f"── 證券交易稅（0.3%）──",
                f"賣出金額：{sell_amount:,.0f} 元",
                f"證交稅：{securities_tax:,.0f} 元",
            ]
        return "\n".join(lines)
    except Exception as e:
        return f"稅務試算失敗：{e}"


def fetch_currency_converter(amount: float, from_currency: str, to_currency: str) -> str:
    try:
        import yfinance as yf
        fc = from_currency.upper()
        tc = to_currency.upper()
        if fc == tc:
            return f"{amount:,.2f} {fc} = {amount:,.2f} {tc}"
        # 嘗試直接查匯率
        pair = f"{fc}{tc}=X"
        ticker = yf.Ticker(pair)
        hist = ticker.history(period="2d")
        if hist.empty:
            # 嘗試反向
            pair2 = f"{tc}{fc}=X"
            hist2 = yf.Ticker(pair2).history(period="2d")
            if hist2.empty:
                return f"無法取得 {fc}/{tc} 匯率資料"
            rate = 1 / hist2["Close"].iloc[-1]
        else:
            rate = hist["Close"].iloc[-1]
        result = amount * rate
        lines = [
            f"💱 外幣換算\n",
            f"{amount:,.2f} {fc}",
            f"= {result:,.4f} {tc}",
            f"\n即時匯率：1 {fc} = {rate:.4f} {tc}",
        ]
        return "\n".join(lines)
    except Exception as e:
        return f"外幣換算失敗：{e}"


def fetch_fund(symbol: str) -> str:
    try:
        import yfinance as yf
        ticker = yf.Ticker(symbol)
        info = ticker.info
        hist = ticker.history(period="1y")
        if hist.empty:
            return f"找不到 {symbol} 的基金資料"
        price = hist["Close"].iloc[-1]
        price_1m = hist["Close"].iloc[-22] if len(hist) > 22 else hist["Close"].iloc[0]
        price_3m = hist["Close"].iloc[-66] if len(hist) > 66 else hist["Close"].iloc[0]
        price_1y = hist["Close"].iloc[0]
        ret_1m = (price / price_1m - 1) * 100
        ret_3m = (price / price_3m - 1) * 100
        ret_1y = (price / price_1y - 1) * 100
        name = info.get("longName") or info.get("shortName") or symbol
        expense = info.get("annualReportExpenseRatio")
        category = info.get("category") or info.get("fundFamily") or "—"
        nav = info.get("navPrice") or price
        lines = [
            f"📦 基金查詢：{name}\n",
            f"代號：{symbol}　類別：{category}",
            f"淨值/價格：{nav:.2f}",
        ]
        if expense:
            lines.append(f"費用率：{expense*100:.2f}%")
        lines += [
            f"",
            f"── 績效 ──",
            f"近1月：{ret_1m:+.2f}%",
            f"近3月：{ret_3m:+.2f}%",
            f"近1年：{ret_1y:+.2f}%",
        ]
        return "\n".join(lines)
    except Exception as e:
        return f"基金查詢失敗：{e}"


def fetch_reits(symbol: str) -> str:
    try:
        import yfinance as yf
        ticker = yf.Ticker(symbol)
        info = ticker.info
        hist = ticker.history(period="1y")
        if hist.empty:
            return f"找不到 {symbol} 的REITs資料"
        price = hist["Close"].iloc[-1]
        price_1y = hist["Close"].iloc[0]
        ret_1y = (price / price_1y - 1) * 100
        name = info.get("longName") or info.get("shortName") or symbol
        div_yield = info.get("dividendYield") or 0
        trailing_annual_div = info.get("trailingAnnualDividendYield") or div_yield
        market_cap = info.get("marketCap") or 0
        sector = info.get("sector") or info.get("category") or "不動產"
        nav = info.get("bookValue") or info.get("navPrice") or 0
        lines = [
            f"🏢 REITs查詢：{name}\n",
            f"代號：{symbol}　類別：{sector}",
            f"現價：{price:.2f}",
        ]
        if div_yield:
            lines.append(f"股息殖利率：{trailing_annual_div*100:.2f}%")
        if market_cap:
            lines.append(f"市值：{market_cap/1e8:.1f} 億")
        if nav:
            premium = (price / nav - 1) * 100
            lines.append(f"NAV：{nav:.2f}（折溢價：{premium:+.1f}%）")
        lines += [
            f"",
            f"近1年報酬：{ret_1y:+.2f}%",
        ]
        return "\n".join(lines)
    except Exception as e:
        return f"REITs查詢失敗：{e}"


def fetch_inflation_adjusted(nominal_return: float, years: int, amount: float,
                              inflation_rate: float = 2.0) -> str:
    try:
        # 費雪方程式：實質報酬 ≈ 名目報酬 - 通膨率
        real_return = ((1 + nominal_return / 100) / (1 + inflation_rate / 100) - 1) * 100
        # 名目終值
        nominal_fv = amount * (1 + nominal_return / 100) ** years
        # 實質終值（通膨調整後的購買力）
        real_fv = amount * (1 + real_return / 100) ** years
        # 通膨吃掉的部分
        inflation_loss = nominal_fv - real_fv
        lines = [
            f"📉 通膨調整報酬試算\n",
            f"本金：{amount:,.0f} 元",
            f"名目報酬率：{nominal_return}%　通膨率：{inflation_rate}%",
            f"實質報酬率：{real_return:.2f}%　期間：{years} 年",
            f"",
            f"── {years}年後 ──",
            f"名目終值：{nominal_fv:,.0f} 元",
            f"實質購買力：{real_fv:,.0f} 元",
            f"通膨吃掉：{inflation_loss:,.0f} 元（{inflation_loss/nominal_fv*100:.1f}%）",
            f"",
            f"💡 今天 {amount:,.0f} 元的東西，{years}年後需要 {amount*(1+inflation_rate/100)**years:,.0f} 元才買得到",
        ]
        return "\n".join(lines)
    except Exception as e:
        return f"通膨調整試算失敗：{e}"


def fetch_defi_calculator(principal_usd: float, apy: float, days: int,
                           compound: bool = True, protocol: str = "") -> str:
    try:
        import yfinance as yf
        # 查詢USD/TWD匯率
        try:
            usdtwd = yf.Ticker("USDTWD=X").history(period="2d")["Close"].iloc[-1]
        except Exception:
            usdtwd = 32.0
        if compound:
            daily_rate = apy / 100 / 365
            final_usd = principal_usd * (1 + daily_rate) ** days
        else:
            final_usd = principal_usd * (1 + apy / 100 * days / 365)
        profit_usd = final_usd - principal_usd
        lines = [
            f"🔗 DeFi收益試算\n",
        ]
        if protocol:
            lines.append(f"協議：{protocol}")
        lines += [
            f"本金：${principal_usd:,.2f} USD（約 {principal_usd*usdtwd:,.0f} TWD）",
            f"APY：{apy}%　質押天數：{days} 天",
            f"複利：{'是' if compound else '否'}",
            f"",
            f"── 試算結果 ──",
            f"到期本利和：${final_usd:,.4f} USD",
            f"獲利：${profit_usd:,.4f} USD（約 {profit_usd*usdtwd:,.0f} TWD）",
            f"實際年化（{days}天）：{(final_usd/principal_usd-1)*365/days*100:.2f}%",
            f"",
            f"⚠️ DeFi有智能合約風險、無常損失風險，本試算僅供參考",
        ]
        return "\n".join(lines)
    except Exception as e:
        return f"DeFi試算失敗：{e}"


def fetch_gold_calculator(weight: float, unit: str = "公克", currency: str = "TWD") -> str:
    try:
        import yfinance as yf
        # 查詢黃金價格（美元/盎司）
        gold = yf.Ticker("GC=F").history(period="2d")
        if gold.empty:
            gold = yf.Ticker("GLD").history(period="2d")
            if gold.empty:
                return "無法取得黃金價格"
            gold_usd_oz = gold["Close"].iloc[-1] / 0.0965835  # GLD每股≈0.0965835盎司
        else:
            gold_usd_oz = gold["Close"].iloc[-1]
        # 換算匯率
        try:
            usdtwd = yf.Ticker("USDTWD=X").history(period="2d")["Close"].iloc[-1]
        except Exception:
            usdtwd = 32.0
        # 單位換算為公克
        unit_map = {"公克": 1, "錢": 3.75, "兩": 37.5, "盎司": 31.1035}
        gram = weight * unit_map.get(unit, 1)
        # 公克換算盎司
        oz = gram / 31.1035
        value_usd = oz * gold_usd_oz
        value_twd = value_usd * usdtwd
        gold_per_gram_twd = gold_usd_oz / 31.1035 * usdtwd
        lines = [
            f"🥇 黃金換算\n",
            f"國際金價：${gold_usd_oz:,.2f}/盎司",
            f"每公克：{gold_per_gram_twd:,.0f} TWD　${gold_usd_oz/31.1035:.2f} USD",
            f"",
            f"── 換算結果 ──",
            f"{weight} {unit} = {gram:.4f} 公克 = {oz:.6f} 盎司",
        ]
        if currency == "TWD":
            lines.append(f"價值：{value_twd:,.0f} 新台幣（匯率 {usdtwd:.2f}）")
        else:
            lines.append(f"價值：${value_usd:,.2f} USD")
        return "\n".join(lines)
    except Exception as e:
        return f"黃金換算失敗：{e}"


def fetch_forex_deposit(amount_twd: float, currency: str, annual_rate: float, months: int,
                         buy_rate: float = None, sell_rate: float = None) -> str:
    try:
        import yfinance as yf
        cur = currency.upper()
        # 查詢即時匯率
        if not buy_rate:
            pair = f"TWD{cur}=X"
            hist = yf.Ticker(pair).history(period="2d")
            if hist.empty:
                pair2 = f"{cur}TWD=X"
                hist2 = yf.Ticker(pair2).history(period="2d")
                if not hist2.empty:
                    rate_per_twd = 1 / hist2["Close"].iloc[-1]
                else:
                    return f"無法取得 {cur}/TWD 匯率"
            else:
                rate_per_twd = hist["Close"].iloc[-1]
            buy_rate_actual = 1 / rate_per_twd  # TWD→外幣
        else:
            rate_per_twd = 1 / buy_rate
            buy_rate_actual = buy_rate
        if not sell_rate:
            sell_rate_actual = buy_rate_actual  # 保守假設相同
        else:
            sell_rate_actual = sell_rate
        # 換算外幣本金
        foreign_principal = amount_twd / buy_rate_actual
        # 計算外幣到期本利和
        foreign_final = foreign_principal * (1 + annual_rate / 100 * months / 12)
        foreign_interest = foreign_final - foreign_principal
        # 換回台幣
        twd_final = foreign_final * sell_rate_actual
        twd_profit = twd_final - amount_twd
        effective_rate = (twd_profit / amount_twd / (months / 12) * 100) if amount_twd > 0 and months > 0 else 0
        lines = [
            f"🌐 外幣定存試算（{cur}）\n",
            f"台幣本金：{amount_twd:,.0f} 元",
            f"買入匯率：1 {cur} = {buy_rate_actual:.4f} TWD",
            f"外幣本金：{foreign_principal:,.2f} {cur}",
            f"年利率：{annual_rate}%　存款期：{months} 個月",
            f"",
            f"── 到期結果 ──",
            f"外幣本利和：{foreign_final:,.4f} {cur}",
            f"外幣利息：{foreign_interest:,.4f} {cur}",
            f"賣出匯率：1 {cur} = {sell_rate_actual:.4f} TWD",
            f"換回台幣：{twd_final:,.0f} 元",
            f"台幣獲利：{twd_profit:+,.0f} 元",
            f"等效台幣年利率：{effective_rate:.2f}%",
        ]
        if abs(buy_rate_actual - sell_rate_actual) < 0.001:
            lines.append("\n⚠️ 未考慮匯差手續費及匯率變動風險")
        return "\n".join(lines)
    except Exception as e:
        return f"外幣定存試算失敗：{e}"


def fetch_financial_health(monthly_income: float, monthly_expense: float,
                            total_assets: float, total_debt: float,
                            emergency_fund_months: float = 0,
                            has_insurance: bool = False,
                            investment_ratio: float = 0) -> str:
    try:
        score = 100
        issues = []
        goods = []
        # 儲蓄率
        save_rate = (monthly_income - monthly_expense) / monthly_income * 100 if monthly_income > 0 else 0
        if save_rate < 0:
            score -= 30; issues.append("每月支出超過收入（負儲蓄）")
        elif save_rate < 10:
            score -= 15; issues.append(f"儲蓄率偏低（{save_rate:.1f}%，建議≥20%）")
        elif save_rate >= 20:
            goods.append(f"儲蓄率良好（{save_rate:.1f}%）")
        # 負債比
        debt_ratio = total_debt / total_assets * 100 if total_assets > 0 else 100
        if debt_ratio > 70:
            score -= 25; issues.append(f"負債比過高（{debt_ratio:.1f}%，建議＜50%）")
        elif debt_ratio > 50:
            score -= 10; issues.append(f"負債比偏高（{debt_ratio:.1f}%）")
        else:
            goods.append(f"負債比健康（{debt_ratio:.1f}%）")
        # 緊急備用金
        if emergency_fund_months < 3:
            score -= 20; issues.append(f"緊急備用金不足（{emergency_fund_months}個月，建議≥6個月）")
        elif emergency_fund_months >= 6:
            goods.append(f"緊急備用金充足（{emergency_fund_months}個月）")
        # 保險
        if not has_insurance:
            score -= 10; issues.append("缺乏壽險/重疾險保障")
        else:
            goods.append("有保險保障")
        # 投資比例
        if investment_ratio >= 20:
            goods.append(f"積極投資（{investment_ratio}%收入）")
        elif investment_ratio > 0:
            issues.append(f"投資比例偏低（{investment_ratio}%，建議≥20%）")
            score -= 5
        # 評分
        score = max(0, min(100, score))
        if score >= 80:
            grade, emoji = "優良", "🟢"
        elif score >= 60:
            grade, emoji = "尚可", "🟡"
        elif score >= 40:
            grade, emoji = "需改善", "🟠"
        else:
            grade, emoji = "高風險", "🔴"
        lines = [
            f"💊 財務健康診斷\n",
            f"月收入：{monthly_income:,.0f}　月支出：{monthly_expense:,.0f}",
            f"總資產：{total_assets/10000:,.0f}萬　總負債：{total_debt/10000:,.0f}萬",
            f"",
            f"── 健康評分 ──",
            f"{emoji} {score} 分 / 100（{grade}）",
            f"儲蓄率：{save_rate:.1f}%　負債比：{debt_ratio:.1f}%",
        ]
        if goods:
            lines.append("\n✅ 優點：" + "、".join(goods))
        if issues:
            lines.append("\n⚠️ 待改善：")
            for i in issues:
                lines.append(f"  • {i}")
        return "\n".join(lines)
    except Exception as e:
        return f"財務健康診斷失敗：{e}"


def fetch_deep_research(topic: str, lang: str = "zh-tw", depth: int = 5) -> str:
    try:
        from ddgs import DDGS
        depth = int(depth) if not isinstance(depth, int) else depth
        depth = min(max(depth, 3), 8)
        # 自動生成子問題
        sub_questions = [
            f"{topic} 是什麼 基本介紹",
            f"{topic} 最新發展 現況",
            f"{topic} 數據 統計 報告",
            f"{topic} 爭議 問題 缺點",
            f"{topic} 未來趨勢 預測",
            f"{topic} 專家看法 分析",
            f"{topic} 影響 重要性",
            f"{topic} 解決方案 建議",
        ][:depth]
        results = {}
        with DDGS() as ddgs:
            for q in sub_questions:
                hits = list(ddgs.text(q, region="tw-tzh" if lang == "zh-tw" else "us-en", max_results=3))
                if hits:
                    results[q] = hits
        lines = [f"📚 深度研究：{topic}\n"]
        for q, hits in results.items():
            lines.append(f"【{q}】")
            for h in hits:
                title = h.get("title", "")
                body = h.get("body", "")[:200]
                lines.append(f"  • {title}：{body}")
            lines.append("")
        lines.append(f"共蒐集 {len(results)} 個面向，{sum(len(v) for v in results.values())} 筆資料")
        return "\n".join(lines)
    except Exception as e:
        return f"深度研究失敗：{e}"


def fetch_fact_check(claim: str, lang: str = "zh-tw") -> str:
    try:
        from ddgs import DDGS
        region = "tw-tzh" if lang == "zh-tw" else "us-en"
        queries = [
            f"{claim} 真假 查核",
            f"{claim} 事實查核",
            f"{claim} 錯誤 謠言",
            f"{claim} 正確 證實",
        ]
        all_hits = []
        with DDGS() as ddgs:
            for q in queries:
                hits = list(ddgs.text(q, region=region, max_results=3))
                all_hits.extend(hits)
        # 簡單關鍵字判斷情緒
        text_blob = " ".join(h.get("title","") + " " + h.get("body","") for h in all_hits).lower()
        false_kw = ["假", "謠言", "錯誤", "誤導", "不實", "false", "fake", "wrong", "misleading"]
        true_kw = ["真", "確認", "屬實", "正確", "true", "confirmed", "correct", "verified"]
        false_score = sum(text_blob.count(k) for k in false_kw)
        true_score = sum(text_blob.count(k) for k in true_kw)
        if false_score > true_score * 2:
            verdict = "❌ 可能為假/誤導"
        elif true_score > false_score * 2:
            verdict = "✅ 可能屬實"
        elif false_score > 0 or true_score > 0:
            verdict = "⚠️ 有爭議，需進一步確認"
        else:
            verdict = "❓ 資料不足，無法判斷"
        lines = [f"🔍 事實查核\n", f"說法：「{claim}」\n", f"查核結果：{verdict}\n", f"── 相關資料 ──"]
        for h in all_hits[:6]:
            title = h.get("title", "")
            body = h.get("body", "")[:150]
            url = h.get("href", "")
            lines.append(f"• {title}\n  {body}")
        return "\n".join(lines)
    except Exception as e:
        return f"事實查核失敗：{e}"


def fetch_timeline_events(topic: str, lang: str = "zh-tw") -> str:
    try:
        from ddgs import DDGS
        region = "tw-tzh" if lang == "zh-tw" else "us-en"
        queries = [
            f"{topic} 歷史 時間軸 發展",
            f"{topic} 大事記 年表",
            f"{topic} 起源 始末 過程",
        ]
        all_hits = []
        with DDGS() as ddgs:
            for q in queries:
                hits = list(ddgs.text(q, region=region, max_results=4))
                all_hits.extend(hits)
        lines = [f"📅 時間軸：{topic}\n"]
        seen = set()
        for h in all_hits:
            title = h.get("title", "")
            if title in seen:
                continue
            seen.add(title)
            body = h.get("body", "")[:300]
            lines.append(f"【{title}】\n{body}\n")
        lines.append(f"（共 {len(seen)} 筆資料，建議搭配 Wikipedia 查詢完整年表）")
        return "\n".join(lines)
    except Exception as e:
        return f"時間軸整理失敗：{e}"


def fetch_sentiment_scan(topic: str, lang: str = "zh-tw") -> str:
    try:
        from ddgs import DDGS
        region = "tw-tzh" if lang == "zh-tw" else "us-en"
        queries = [
            f"{topic} 正面 支持 優點",
            f"{topic} 負面 反對 批評 缺點",
            f"{topic} 民眾看法 輿論 評價",
        ]
        pos_hits, neg_hits, neutral_hits = [], [], []
        with DDGS() as ddgs:
            pos_hits = list(ddgs.text(queries[0], region=region, max_results=4))
            neg_hits = list(ddgs.text(queries[1], region=region, max_results=4))
            neutral_hits = list(ddgs.text(queries[2], region=region, max_results=3))
        total = len(pos_hits) + len(neg_hits) + len(neutral_hits)
        pos_pct = round(len(pos_hits) / total * 100) if total else 0
        neg_pct = round(len(neg_hits) / total * 100) if total else 0
        neu_pct = 100 - pos_pct - neg_pct
        lines = [
            f"📊 輿情掃描：{topic}\n",
            f"正面 {pos_pct}% ｜ 負面 {neg_pct}% ｜ 中立 {neu_pct}%\n",
            f"── 正面觀點 ──",
        ]
        for h in pos_hits[:2]:
            lines.append(f"• {h.get('title','')}：{h.get('body','')[:120]}")
        lines.append(f"\n── 負面觀點 ──")
        for h in neg_hits[:2]:
            lines.append(f"• {h.get('title','')}：{h.get('body','')[:120]}")
        lines.append(f"\n── 中立/綜合 ──")
        for h in neutral_hits[:2]:
            lines.append(f"• {h.get('title','')}：{h.get('body','')[:120]}")
        return "\n".join(lines)
    except Exception as e:
        return f"輿情掃描失敗：{e}"


def fetch_compare_analysis(items: list, dimensions: list = None, context: str = "") -> str:
    try:
        from ddgs import DDGS
        ctx = f" {context}" if context else ""
        data = {}
        with DDGS() as ddgs:
            for item in items[:5]:
                hits = list(ddgs.text(f"{item}{ctx} 評價 特點 優缺點", region="tw-tzh", max_results=3))
                data[item] = " ".join(h.get("body","") for h in hits)[:500]
        lines = [f"⚖️ 比較分析：{' vs '.join(items)}\n"]
        if context:
            lines.append(f"背景：{context}\n")
        for item in items[:5]:
            lines.append(f"【{item}】")
            lines.append(data.get(item, "資料不足")[:400])
            lines.append("")
        lines.append("── 綜合建議 ──")
        lines.append(f"以上為各項資料彙整，請根據您的需求與優先考量做最終選擇。")
        return "\n".join(lines)
    except Exception as e:
        return f"比較分析失敗：{e}"


def fetch_pros_cons_analysis(subject: str, context: str = "", lang: str = "zh-tw") -> str:
    try:
        from ddgs import DDGS
        region = "tw-tzh" if lang == "zh-tw" else "us-en"
        ctx = f" {context}" if context else ""
        pros_hits, cons_hits = [], []
        with DDGS() as ddgs:
            pros_hits = list(ddgs.text(f"{subject}{ctx} 優點 好處 支持", region=region, max_results=4))
            cons_hits = list(ddgs.text(f"{subject}{ctx} 缺點 壞處 風險 問題", region=region, max_results=4))
        lines = [f"📋 優缺點分析：{subject}\n"]
        if context:
            lines.append(f"背景：{context}\n")
        lines.append("── 優點 / 支持論點 ──")
        for h in pros_hits[:3]:
            lines.append(f"✅ {h.get('title','')}：{h.get('body','')[:150]}")
        lines.append("\n── 缺點 / 反對論點 ──")
        for h in cons_hits[:3]:
            lines.append(f"⚠️ {h.get('title','')}：{h.get('body','')[:150]}")
        confidence = "高" if len(pros_hits) + len(cons_hits) >= 6 else "中"
        lines.append(f"\n資料信心度：{confidence}（共 {len(pros_hits)+len(cons_hits)} 筆）")
        return "\n".join(lines)
    except Exception as e:
        return f"優缺點分析失敗：{e}"


def fetch_research_report(topic: str, purpose: str = "一般研究", lang: str = "zh-tw") -> str:
    try:
        from ddgs import DDGS
        region = "tw-tzh" if lang == "zh-tw" else "us-en"
        sections = {
            "背景與定義": f"{topic} 定義 介紹 背景",
            "現況與數據": f"{topic} 現況 統計 數據 規模",
            "主要發現": f"{topic} 研究 發現 結果 報告",
            "爭議與挑戰": f"{topic} 問題 挑戰 爭議",
            "趨勢與展望": f"{topic} 趨勢 未來 預測 展望",
        }
        collected = {}
        with DDGS() as ddgs:
            for sec, q in sections.items():
                hits = list(ddgs.text(q, region=region, max_results=3))
                collected[sec] = hits
        lines = [
            f"📄 研究報告：{topic}",
            f"目的：{purpose}\n",
            f"═══ 執行摘要 ═══",
            f"本報告針對「{topic}」進行多面向資料蒐集，涵蓋背景、數據、爭議與展望。\n",
        ]
        for sec, hits in collected.items():
            lines.append(f"── {sec} ──")
            for h in hits[:2]:
                lines.append(f"• {h.get('title','')}：{h.get('body','')[:200]}")
            lines.append("")
        lines += [
            "═══ 結論與建議 ═══",
            f"根據蒐集資料，「{topic}」是一個值得深入關注的議題。",
            f"建議進一步參閱原始來源以獲得更完整資訊。",
        ]
        return "\n".join(lines)
    except Exception as e:
        return f"研究報告生成失敗：{e}"


def fetch_opinion_writer(topic: str, stance: str = "中立", style: str = "正式") -> str:
    try:
        from ddgs import DDGS
        hits = []
        with DDGS() as ddgs:
            hits = list(ddgs.text(f"{topic} 分析 觀點 評論", region="tw-tzh", max_results=6))
        context = "\n".join(f"- {h.get('title','')}：{h.get('body','')[:200]}" for h in hits)
        lines = [f"✍️ 觀點撰寫：{topic}",
                 f"立場：{stance}　文風：{style}\n",
                 f"── 資料基礎 ──",
                 context,
                 f"\n── {stance}立場分析 ──",
        ]
        if stance == "支持":
            lines.append(f"從現有資料來看，「{topic}」有其正面價值。以下論點支持此立場：")
        elif stance == "反對":
            lines.append(f"從現有資料來看，「{topic}」存在值得警惕的問題。以下論點提出質疑：")
        elif stance == "批判":
            lines.append(f"以批判性視角審視「{topic}」，可發現以下值得深究之處：")
        else:
            lines.append(f"綜合現有資料，「{topic}」可從多角度理解：")
        lines.append(f"（本節由 Claude 依蒐集資料整合後發表看法）")
        return "\n".join(lines)
    except Exception as e:
        return f"觀點撰寫失敗：{e}"


def fetch_trend_forecast(topic: str, timeframe: str = "全部", lang: str = "zh-tw") -> str:
    try:
        from ddgs import DDGS
        region = "tw-tzh" if lang == "zh-tw" else "us-en"
        queries = [
            f"{topic} 未來趨勢 預測 展望",
            f"{topic} 短期 發展 2024 2025",
            f"{topic} 長期 影響 趨勢",
        ]
        all_hits = []
        with DDGS() as ddgs:
            for q in queries:
                hits = list(ddgs.text(q, region=region, max_results=3))
                all_hits.extend(hits)
        lines = [f"🔮 趨勢預測：{topic}", f"預測範圍：{timeframe}\n"]
        if timeframe in ("短期(1年內)", "全部"):
            lines.append("── 短期（1年內）──")
            for h in all_hits[:3]:
                lines.append(f"• {h.get('title','')}：{h.get('body','')[:180]}")
            lines.append("")
        if timeframe in ("中期(1-3年)", "全部"):
            lines.append("── 中期（1–3年）──")
            for h in all_hits[3:6]:
                lines.append(f"• {h.get('title','')}：{h.get('body','')[:180]}")
            lines.append("")
        if timeframe in ("長期(3年以上)", "全部"):
            lines.append("── 長期（3年以上）──")
            for h in all_hits[6:9]:
                lines.append(f"• {h.get('title','')}：{h.get('body','')[:180]}")
            lines.append("")
        lines.append(f"⚠️ 預測基於現有公開資料，實際發展受多重因素影響")
        return "\n".join(lines)
    except Exception as e:
        return f"趨勢預測失敗：{e}"


def fetch_debate_simulator(motion: str, lang: str = "zh-tw") -> str:
    try:
        from ddgs import DDGS
        region = "tw-tzh" if lang == "zh-tw" else "us-en"
        pro_hits, con_hits = [], []
        with DDGS() as ddgs:
            pro_hits = list(ddgs.text(f"{motion} 支持 贊成 優點 好處", region=region, max_results=4))
            con_hits = list(ddgs.text(f"{motion} 反對 質疑 缺點 問題", region=region, max_results=4))
        lines = [
            f"⚔️ 辯論模擬：{motion}\n",
            f"══ 正方論點 ══",
        ]
        for h in pro_hits[:3]:
            lines.append(f"✅ {h.get('title','')}：{h.get('body','')[:180]}")
        lines += [f"\n══ 反方論點 ══"]
        for h in con_hits[:3]:
            lines.append(f"❌ {h.get('title','')}：{h.get('body','')[:180]}")
        pro_strength = len(pro_hits)
        con_strength = len(con_hits)
        if pro_strength > con_strength:
            verdict = "正方論據較充分"
        elif con_strength > pro_strength:
            verdict = "反方論據較充分"
        else:
            verdict = "雙方論據相當，難以定論"
        lines += [
            f"\n══ 綜合判斷 ══",
            f"議題：「{motion}」",
            f"評估：{verdict}",
            f"此議題涉及多方面考量，建議結合個人價值觀與具體情境做判斷。",
        ]
        return "\n".join(lines)
    except Exception as e:
        return f"辯論模擬失敗：{e}"


def fetch_academic_search(query: str, field: str = "", lang: str = "en") -> str:
    try:
        from ddgs import DDGS
        region = "us-en" if lang == "en" else "tw-tzh"
        field_tag = f" {field}" if field else ""
        queries = [
            f"{query}{field_tag} site:scholar.google.com OR site:pubmed.ncbi.nlm.nih.gov OR site:semanticscholar.org",
            f"{query}{field_tag} research study findings",
            f"{query}{field_tag} academic paper review",
        ]
        all_hits = []
        with DDGS() as ddgs:
            for q in queries:
                hits = list(ddgs.text(q, region=region, max_results=4))
                all_hits.extend(hits)
        seen, lines = set(), [f"🎓 學術搜尋：{query}\n"]
        if field:
            lines.append(f"領域：{field}\n")
        count = 0
        for h in all_hits:
            title = h.get("title", "")
            if title in seen:
                continue
            seen.add(title)
            body = h.get("body", "")[:250]
            url = h.get("href", "")
            lines.append(f"📄 {title}\n   {body}")
            if url:
                lines.append(f"   {url}")
            lines.append("")
            count += 1
            if count >= 6:
                break
        lines.append(f"共找到 {count} 篇相關學術資料（建議至 Google Scholar 完整查閱）")
        return "\n".join(lines)
    except Exception as e:
        return f"學術搜尋失敗：{e}"


def fetch_health_research(topic: str, lang: str = "zh-tw") -> str:
    try:
        from ddgs import DDGS
        region = "tw-tzh" if lang == "zh-tw" else "us-en"
        queries = [
            f"{topic} 症狀 原因 說明",
            f"{topic} 治療 建議 注意事項",
            f"{topic} 衛福部 醫學 研究",
        ]
        all_hits = []
        with DDGS() as ddgs:
            for q in queries:
                hits = list(ddgs.text(q, region=region, max_results=3))
                all_hits.extend(hits)
        lines = [
            f"🏥 健康資訊：{topic}\n",
            f"⚠️ 以下資訊僅供參考，不替代醫師診斷，如有不適請就醫。\n",
        ]
        seen = set()
        for h in all_hits[:6]:
            title = h.get("title", "")
            if title in seen:
                continue
            seen.add(title)
            body = h.get("body", "")[:250]
            lines.append(f"• {title}\n  {body}\n")
        return "\n".join(lines)
    except Exception as e:
        return f"健康資訊搜尋失敗：{e}"


def fetch_law_research(topic: str, jurisdiction: str = "台灣", lang: str = "zh-tw") -> str:
    try:
        from ddgs import DDGS
        region = "tw-tzh" if lang == "zh-tw" else "us-en"
        queries = [
            f"{jurisdiction} {topic} 法律 法規 條文",
            f"{jurisdiction} {topic} 判例 實務 見解",
            f"{topic} 法律問題 解答",
        ]
        all_hits = []
        with DDGS() as ddgs:
            for q in queries:
                hits = list(ddgs.text(q, region=region, max_results=3))
                all_hits.extend(hits)
        lines = [
            f"⚖️ 法規查詢：{topic}（{jurisdiction}）\n",
            f"⚠️ 以下資訊僅供參考，不構成法律意見，具體情況建議諮詢律師。\n",
        ]
        seen = set()
        for h in all_hits[:6]:
            title = h.get("title", "")
            if title in seen:
                continue
            seen.add(title)
            body = h.get("body", "")[:250]
            lines.append(f"• {title}\n  {body}\n")
        return "\n".join(lines)
    except Exception as e:
        return f"法規查詢失敗：{e}"


def fetch_person_research(name: str, context: str = "", lang: str = "zh-tw") -> str:
    try:
        from ddgs import DDGS
        region = "tw-tzh" if lang == "zh-tw" else "us-en"
        ctx = f" {context}" if context else ""
        queries = [
            f"{name}{ctx} 背景 經歷 介紹",
            f"{name}{ctx} 成就 評價 貢獻",
            f"{name}{ctx} 爭議 批評 問題",
        ]
        sections = {"背景與經歷": [], "成就與評價": [], "爭議與批評": []}
        sec_keys = list(sections.keys())
        with DDGS() as ddgs:
            for i, q in enumerate(queries):
                hits = list(ddgs.text(q, region=region, max_results=3))
                sections[sec_keys[i]] = hits
        lines = [f"👤 人物研究：{name}\n"]
        if context:
            lines.append(f"背景：{context}\n")
        for sec, hits in sections.items():
            if hits:
                lines.append(f"── {sec} ──")
                for h in hits[:2]:
                    lines.append(f"• {h.get('title','')}：{h.get('body','')[:200]}")
                lines.append("")
        return "\n".join(lines)
    except Exception as e:
        return f"人物研究失敗：{e}"


def fetch_company_research(company: str, lang: str = "zh-tw") -> str:
    try:
        from ddgs import DDGS
        import yfinance as yf
        region = "tw-tzh" if lang == "zh-tw" else "us-en"
        lines = [f"🏢 公司研究：{company}\n"]
        # 嘗試取得財務數據
        try:
            ticker = yf.Ticker(company)
            info = ticker.info
            name = info.get("longName") or info.get("shortName") or company
            sector = info.get("sector") or "—"
            industry = info.get("industry") or "—"
            employees = info.get("fullTimeEmployees") or "—"
            revenue = info.get("totalRevenue")
            market_cap = info.get("marketCap")
            lines += [
                f"── 基本資料 ──",
                f"公司：{name}　產業：{sector} / {industry}",
            ]
            if employees != "—":
                lines.append(f"員工數：{employees:,}")
            if revenue:
                lines.append(f"年營收：${revenue/1e8:.1f}億")
            if market_cap:
                lines.append(f"市值：${market_cap/1e8:.1f}億")
            lines.append("")
        except Exception:
            pass
        # 搜尋新聞與評價
        queries = {
            "公司動態": f"{company} 最新消息 發展",
            "產品評價": f"{company} 產品 服務 評價",
            "競爭分析": f"{company} 競爭對手 市場地位",
        }
        with DDGS() as ddgs:
            for sec, q in queries.items():
                hits = list(ddgs.text(q, region=region, max_results=2))
                if hits:
                    lines.append(f"── {sec} ──")
                    for h in hits:
                        lines.append(f"• {h.get('title','')}：{h.get('body','')[:180]}")
                    lines.append("")
        return "\n".join(lines)
    except Exception as e:
        return f"公司研究失敗：{e}"


def fetch_product_review(product: str, category: str = "", lang: str = "zh-tw") -> str:
    try:
        from ddgs import DDGS
        region = "tw-tzh" if lang == "zh-tw" else "us-en"
        cat = f" {category}" if category else ""
        queries = [
            f"{product}{cat} 評測 開箱 評價",
            f"{product}{cat} 優點 缺點 推薦",
            f"{product}{cat} 使用心得 評分",
        ]
        all_hits = []
        with DDGS() as ddgs:
            for q in queries:
                hits = list(ddgs.text(q, region=region, max_results=3))
                all_hits.extend(hits)
        pos_kw = ["推薦", "好用", "優秀", "值得", "滿意", "excellent", "great", "recommend"]
        neg_kw = ["不推", "缺點", "問題", "失望", "差", "poor", "bad", "issue"]
        all_text = " ".join(h.get("body","") for h in all_hits).lower()
        pos_score = sum(all_text.count(k) for k in pos_kw)
        neg_score = sum(all_text.count(k) for k in neg_kw)
        total = pos_score + neg_score
        rating = round(pos_score / total * 5, 1) if total > 0 else 3.0
        lines = [
            f"⭐ 產品評測：{product}",
            f"綜合評分：{rating}/5.0　（正面{pos_score}則 / 負面{neg_score}則）\n",
            f"── 評測彙整 ──",
        ]
        seen = set()
        for h in all_hits[:6]:
            title = h.get("title","")
            if title in seen:
                continue
            seen.add(title)
            body = h.get("body","")[:200]
            lines.append(f"• {title}：{body}\n")
        return "\n".join(lines)
    except Exception as e:
        return f"產品評測失敗：{e}"


def fetch_travel_research(destination: str, days: int = None, style: str = "", lang: str = "zh-tw") -> str:
    try:
        from ddgs import DDGS
        region = "tw-tzh" if lang == "zh-tw" else "us-en"
        style_tag = f" {style}" if style else ""
        days_tag = f" {days}天" if days else ""
        queries = {
            "景點": f"{destination} 必去景點 推薦",
            "美食": f"{destination} 必吃美食 餐廳",
            "交通住宿": f"{destination} 交通 住宿 費用",
            "注意事項": f"{destination} 旅遊注意 簽證 安全",
        }
        lines = [f"✈️ 旅遊研究：{destination}"]
        if days:
            lines.append(f"行程天數：{days} 天")
        if style:
            lines.append(f"旅遊風格：{style}")
        lines.append("")
        with DDGS() as ddgs:
            for sec, q in queries.items():
                hits = list(ddgs.text(q + style_tag + days_tag, region=region, max_results=2))
                if hits:
                    lines.append(f"── {sec} ──")
                    for h in hits:
                        lines.append(f"• {h.get('title','')}：{h.get('body','')[:200]}")
                    lines.append("")
        return "\n".join(lines)
    except Exception as e:
        return f"旅遊研究失敗：{e}"


def fetch_job_market(job_title: str, location: str = "台灣", lang: str = "zh-tw") -> str:
    try:
        from ddgs import DDGS
        region = "tw-tzh" if lang == "zh-tw" else "us-en"
        queries = {
            "薪資行情": f"{location} {job_title} 薪資 薪水 行情",
            "技能需求": f"{job_title} 必備技能 技術要求 條件",
            "市場需求": f"{location} {job_title} 職缺 需求 前景",
            "未來趨勢": f"{job_title} 產業趨勢 未來發展 AI影響",
        }
        lines = [f"💼 職涯市場分析：{job_title}（{location}）\n"]
        with DDGS() as ddgs:
            for sec, q in queries.items():
                hits = list(ddgs.text(q, region=region, max_results=2))
                if hits:
                    lines.append(f"── {sec} ──")
                    for h in hits:
                        lines.append(f"• {h.get('title','')}：{h.get('body','')[:200]}")
                    lines.append("")
        return "\n".join(lines)
    except Exception as e:
        return f"職涯市場分析失敗：{e}"


def fetch_impact_analysis(event: str, scope: list = None, lang: str = "zh-tw") -> str:
    try:
        from ddgs import DDGS
        region = "tw-tzh" if lang == "zh-tw" else "us-en"
        if not scope:
            scope = ["個人", "企業", "社會", "經濟"]
        lines = [f"🌐 影響力分析：{event}\n"]
        with DDGS() as ddgs:
            for s in scope:
                hits = list(ddgs.text(f"{event} 對{s}的影響", region=region, max_results=2))
                if hits:
                    lines.append(f"── 對{s}的影響 ──")
                    for h in hits:
                        lines.append(f"• {h.get('title','')}：{h.get('body','')[:200]}")
                    lines.append("")
        return "\n".join(lines)
    except Exception as e:
        return f"影響力分析失敗：{e}"


def fetch_scenario_planning(topic: str, horizon: str = "", lang: str = "zh-tw") -> str:
    try:
        from ddgs import DDGS
        region = "tw-tzh" if lang == "zh-tw" else "us-en"
        h_tag = f" {horizon}" if horizon else ""
        scenarios = {
            "樂觀情境": f"{topic}{h_tag} 最好情況 成功 機會",
            "基準情境": f"{topic}{h_tag} 預測 可能發展 趨勢",
            "悲觀情境": f"{topic}{h_tag} 風險 失敗 最壞情況",
        }
        lines = [f"🔭 情境規劃：{topic}"]
        if horizon:
            lines.append(f"時間範圍：{horizon}")
        lines.append("")
        with DDGS() as ddgs:
            for sc, q in scenarios.items():
                hits = list(ddgs.text(q, region=region, max_results=3))
                lines.append(f"══ {sc} ══")
                for h in hits[:2]:
                    lines.append(f"• {h.get('title','')}：{h.get('body','')[:200]}")
                lines.append("")
        lines.append("💡 建議針對各情境預先制定因應策略，提高應變能力")
        return "\n".join(lines)
    except Exception as e:
        return f"情境規劃失敗：{e}"


def fetch_decision_helper(question: str, options: list = None, criteria: list = None) -> str:
    try:
        from ddgs import DDGS
        lines = [f"🤔 決策輔助：{question}\n"]
        if options:
            lines.append(f"選項：{' vs '.join(options)}\n")
        if criteria:
            lines.append(f"考量標準：{', '.join(criteria)}\n")
        # 搜尋相關資訊
        with DDGS() as ddgs:
            hits = list(ddgs.text(f"{question} 建議 分析 怎麼決定", region="tw-tzh", max_results=5))
        lines.append("── 相關資訊 ──")
        for h in hits[:4]:
            lines.append(f"• {h.get('title','')}：{h.get('body','')[:200]}")
        lines.append("")
        if options:
            lines.append("── 各選項分析 ──")
            with DDGS() as ddgs:
                for opt in options[:3]:
                    opt_hits = list(ddgs.text(f"{opt} 優缺點 評價", region="tw-tzh", max_results=2))
                    lines.append(f"【{opt}】")
                    for h in opt_hits[:1]:
                        lines.append(f"  {h.get('body','')[:180]}")
                    lines.append("")
        lines.append("── 建議框架 ──")
        if criteria:
            for c in criteria:
                lines.append(f"□ {c}：請根據您的具體情況評分")
        lines.append("\n綜合以上資訊，建議依個人優先順序做最終判斷。")
        return "\n".join(lines)
    except Exception as e:
        return f"決策輔助失敗：{e}"


def fetch_devil_advocate(position: str, lang: str = "zh-tw") -> str:
    try:
        from ddgs import DDGS
        region = "tw-tzh" if lang == "zh-tw" else "us-en"
        queries = [
            f"{position} 反對 批評 缺點 問題",
            f"{position} 失敗案例 風險 危險",
            f"{position} 質疑 挑戰 反駁",
        ]
        all_hits = []
        with DDGS() as ddgs:
            for q in queries:
                hits = list(ddgs.text(q, region=region, max_results=3))
                all_hits.extend(hits)
        lines = [
            f"😈 魔鬼代言人：挑戰「{position}」\n",
            f"以下從相反角度提出最強反駁，幫助您找出盲點：\n",
            f"── 反駁論點 ──",
        ]
        seen = set()
        for h in all_hits[:6]:
            title = h.get("title","")
            if title in seen:
                continue
            seen.add(title)
            body = h.get("body","")[:200]
            lines.append(f"⚡ {title}：{body}\n")
        lines.append("── 結語 ──")
        lines.append("以上為刻意的反面論點，目的是強化您的思考。若能反駁以上論點，您的立場將更為穩固。")
        return "\n".join(lines)
    except Exception as e:
        return f"魔鬼代言人失敗：{e}"


def fetch_summary_writer(topic: str, max_points: int = 7, lang: str = "zh-tw") -> str:
    try:
        from ddgs import DDGS
        region = "tw-tzh" if lang == "zh-tw" else "us-en"
        all_hits = []
        with DDGS() as ddgs:
            for q in [topic, f"{topic} 重點 整理", f"{topic} 分析 摘要"]:
                hits = list(ddgs.text(q, region=region, max_results=4))
                all_hits.extend(hits)
        # 去重並提取文字
        seen, texts = set(), []
        for h in all_hits:
            title = h.get("title", "")
            body = h.get("body", "")
            if title not in seen and body:
                seen.add(title)
                texts.append(f"{title}：{body}")
        combined = "\n".join(texts[:10])
        # 從文字中提取句子當重點
        import re
        sentences = re.split(r'[。！？\n]', combined)
        sentences = [s.strip() for s in sentences if len(s.strip()) > 15][:max_points * 3]
        # 去重相似句
        points, seen_short = [], set()
        for s in sentences:
            key = s[:10]
            if key not in seen_short:
                seen_short.add(key)
                points.append(s[:120])
            if len(points) >= max_points:
                break
        lines = [f"📝 摘要：{topic}\n", f"── 核心重點（{len(points)} 項）──"]
        for i, p in enumerate(points, 1):
            lines.append(f"{i}. {p}")
        lines.append(f"\n共整合 {len(seen)} 篇資料來源")
        return "\n".join(lines)
    except Exception as e:
        return f"摘要失敗：{e}"


def fetch_key_insights(topic: str, count: int = 5, lang: str = "zh-tw") -> str:
    try:
        from ddgs import DDGS
        region = "tw-tzh" if lang == "zh-tw" else "us-en"
        queries = [
            f"{topic} 關鍵發現 重要結論",
            f"{topic} 研究結果 數據 證明",
            f"{topic} 專家觀點 核心",
        ]
        all_hits = []
        with DDGS() as ddgs:
            for q in queries:
                hits = list(ddgs.text(q, region=region, max_results=4))
                all_hits.extend(hits)
        # 找含數字或強烈結論的句子（通常是洞察）
        import re
        insight_kw = ["發現", "證明", "研究", "顯示", "指出", "表明", "關鍵", "重要", "核心",
                      "shows", "reveals", "found", "key", "critical", "significant"]
        candidates = []
        for h in all_hits:
            body = h.get("body", "")
            title = h.get("title", "")
            score = sum(body.count(k) + title.count(k) for k in insight_kw)
            has_num = bool(re.search(r'\d+[%倍億萬]', body))
            candidates.append((score + (2 if has_num else 0), title, body[:200]))
        candidates.sort(reverse=True)
        lines = [f"💡 核心洞察：{topic}\n"]
        seen_t = set()
        idx = 1
        for score, title, body in candidates:
            if title in seen_t or not body:
                continue
            seen_t.add(title)
            lines.append(f"#{idx} {title}\n   → {body}\n")
            idx += 1
            if idx > count:
                break
        lines.append(f"洞察來自 {len(all_hits)} 筆資料，依相關性與數據強度排序")
        return "\n".join(lines)
    except Exception as e:
        return f"洞察萃取失敗：{e}"


def fetch_bias_detector(topic: str, lang: str = "zh-tw") -> str:
    try:
        from ddgs import DDGS
        region = "tw-tzh" if lang == "zh-tw" else "us-en"
        queries = [
            f"{topic} 媒體偏見 立場",
            f"{topic} 支持方 反對方 爭議",
            f"{topic} 批評 質疑 偏頗",
        ]
        all_hits = []
        with DDGS() as ddgs:
            for q in queries:
                hits = list(ddgs.text(q, region=region, max_results=3))
                all_hits.extend(hits)
        # 分析立場傾向
        pro_kw = ["支持", "推崇", "認同", "正面", "優秀", "progressive", "support", "advocate"]
        con_kw = ["反對", "批評", "質疑", "保守", "警告", "oppose", "criticize", "concern"]
        left_kw = ["進步", "左派", "liberal", "progressive", "left"]
        right_kw = ["保守", "右派", "conservative", "right", "traditional"]
        all_text = " ".join(h.get("body","") + h.get("title","") for h in all_hits).lower()
        pro_s = sum(all_text.count(k) for k in pro_kw)
        con_s = sum(all_text.count(k) for k in con_kw)
        left_s = sum(all_text.count(k) for k in left_kw)
        right_s = sum(all_text.count(k) for k in right_kw)
        bias_dir = "中立" if abs(pro_s - con_s) < 3 else ("偏正面/支持" if pro_s > con_s else "偏負面/批評")
        political = "中立" if abs(left_s - right_s) < 2 else ("偏進步/左派" if left_s > right_s else "偏保守/右派")
        lines = [
            f"🔍 偏見偵測：{topic}\n",
            f"── 立場分析 ──",
            f"情感傾向：{bias_dir}（正面訊號{pro_s} vs 負面訊號{con_s}）",
            f"政治傾向：{political}",
            f"",
            f"── 各方觀點樣本 ──",
        ]
        seen = set()
        for h in all_hits[:6]:
            title = h.get("title","")
            if title in seen:
                continue
            seen.add(title)
            body = h.get("body","")[:150]
            lines.append(f"• {title}：{body}")
        lines.append(f"\n⚠️ 閱讀此議題資料時建議多方比對，避免單一來源形成片面認知")
        return "\n".join(lines)
    except Exception as e:
        return f"偏見偵測失敗：{e}"


def fetch_second_opinion(question: str, experts: list = None, lang: str = "zh-tw") -> str:
    try:
        from ddgs import DDGS
        region = "tw-tzh" if lang == "zh-tw" else "us-en"
        if not experts:
            experts = ["經濟學家", "心理學家", "社會學家", "科技專家", "實務工作者"]
        lines = [f"🎓 多專家視角：{question}\n"]
        with DDGS() as ddgs:
            for expert in experts[:5]:
                hits = list(ddgs.text(f"{question} {expert} 觀點 看法", region=region, max_results=2))
                lines.append(f"── {expert}的角度 ──")
                if hits:
                    for h in hits[:1]:
                        lines.append(f"• {h.get('title','')}：{h.get('body','')[:200]}")
                else:
                    lines.append(f"• 資料不足，需進一步搜尋")
                lines.append("")
        lines.append("💡 不同專業背景會產生截然不同的分析框架，綜合參考才能形成全面判斷")
        return "\n".join(lines)
    except Exception as e:
        return f"多專家視角失敗：{e}"


def fetch_brainstorm(problem: str, count: int = 8, style: str = "實用", lang: str = "zh-tw") -> str:
    try:
        from ddgs import DDGS
        region = "tw-tzh" if lang == "zh-tw" else "us-en"
        style_map = {
            "實用": "解決方案 方法 做法",
            "創意": "創新 新穎 非傳統 獨特",
            "顛覆": "顛覆 革命性 完全不同 打破慣例",
        }
        style_q = style_map.get(style, "解決方案 方法")
        all_hits = []
        with DDGS() as ddgs:
            for q in [f"{problem} {style_q}", f"{problem} 案例 成功", f"{problem} 創意 想法"]:
                hits = list(ddgs.text(q, region=region, max_results=3))
                all_hits.extend(hits)
        lines = [f"🧠 腦力激盪：{problem}", f"風格：{style}　目標：{count} 個方案\n"]
        # 從搜尋結果萃取方向
        seen, ideas = set(), []
        for h in all_hits:
            title = h.get("title","")
            body = h.get("body","")
            if title not in seen and body:
                seen.add(title)
                ideas.append(f"{title}：{body[:120]}")
        lines.append("── 方案清單 ──")
        for i, idea in enumerate(ideas[:count], 1):
            lines.append(f"{i}. {idea}\n")
        if len(ideas) < count:
            lines.append(f"（已蒐集 {len(ideas)} 個方向，可進一步細化）")
        return "\n".join(lines)
    except Exception as e:
        return f"腦力激盪失敗：{e}"


def fetch_benchmark_analysis(subject: str, industry: str = "", lang: str = "zh-tw") -> str:
    try:
        from ddgs import DDGS
        region = "tw-tzh" if lang == "zh-tw" else "us-en"
        ind_tag = f" {industry}" if industry else ""
        queries = {
            "業界標竿": f"{industry or subject} 最佳實踐 業界標準 領導者",
            "對標案例": f"{subject}{ind_tag} 對標 學習 比較 領先",
            "改進方向": f"{subject}{ind_tag} 改善 優化 提升 差距",
        }
        lines = [f"📐 標竿分析：{subject}"]
        if industry:
            lines.append(f"產業：{industry}")
        lines.append("")
        with DDGS() as ddgs:
            for sec, q in queries.items():
                hits = list(ddgs.text(q, region=region, max_results=3))
                lines.append(f"── {sec} ──")
                for h in hits[:2]:
                    lines.append(f"• {h.get('title','')}：{h.get('body','')[:200]}")
                lines.append("")
        lines.append("── 建議行動 ──")
        lines.append(f"1. 找出領域內 Top 3 標竿對象，深入研究其核心做法")
        lines.append(f"2. 識別 {subject} 與標竿的具體差距")
        lines.append(f"3. 制定可量化的追趕目標與時間表")
        return "\n".join(lines)
    except Exception as e:
        return f"標竿分析失敗：{e}"


def fetch_steel_man(opposing_view: str, own_position: str = "", lang: str = "zh-tw") -> str:
    try:
        from ddgs import DDGS
        region = "tw-tzh" if lang == "zh-tw" else "us-en"
        # 搜尋支持對立觀點的最強論據
        pro_hits, counter_hits = [], []
        with DDGS() as ddgs:
            pro_hits = list(ddgs.text(f"{opposing_view} 支持 論據 最強理由", region=region, max_results=4))
            if own_position:
                counter_hits = list(ddgs.text(f"{own_position} 論據 支持", region=region, max_results=3))
        lines = [
            f"⚔️ 鋼人論證\n",
            f"對立觀點：「{opposing_view}」\n",
            f"══ 鋼人化：對方最強論點 ══",
            f"（以下為對方觀點的最有力版本，非我方立場）\n",
        ]
        for h in pro_hits[:4]:
            lines.append(f"✦ {h.get('title','')}：{h.get('body','')[:180]}\n")
        if own_position:
            lines += [
                f"══ 我方回應 ══",
                f"立場：「{own_position}」\n",
            ]
            for h in counter_hits[:3]:
                lines.append(f"→ {h.get('title','')}：{h.get('body','')[:180]}\n")
        lines.append("💡 鋼人論證要求：先真正理解對方最強版本，才能給出真正有力的回應。")
        return "\n".join(lines)
    except Exception as e:
        return f"鋼人論證失敗：{e}"


def fetch_socratic_questioning(topic: str, depth: int = 5, lang: str = "zh-tw") -> str:
    try:
        from ddgs import DDGS
        region = "tw-tzh" if lang == "zh-tw" else "us-en"
        # 搜尋議題背景，生成有依據的問題層次
        with DDGS() as ddgs:
            hits = list(ddgs.text(f"{topic} 核心問題 爭議 本質", region=region, max_results=4))
        context = " ".join(h.get("body","")[:150] for h in hits)
        # 問題層次設計
        layers = [
            ("釐清概念", f"「{topic}」的核心定義是什麼？我們所說的究竟指的是哪個層面？"),
            ("探究假設", f"這個說法背後有哪些未被檢驗的假設？是否所有人都認同這些前提？"),
            ("檢驗證據", f"支持這個立場的證據有多可靠？有沒有相反的證據被忽略？"),
            ("探索觀點", f"從不同利益關係人的角度來看，這件事會有什麼不同的詮釋？"),
            ("追問影響", f"如果這個觀點是對的，它會帶來什麼後果？我們準備好接受這些後果了嗎？"),
            ("質疑問題本身", f"我們是否問了正確的問題？有沒有更根本的問題應該先被回答？"),
            ("尋找矛盾", f"這個立場內部有沒有自相矛盾之處？邊界條件在哪裡？"),
            ("回歸本質", f"剝除所有表象後，這個議題的最核心本質究竟是什麼？"),
        ][:depth]
        lines = [f"🏛️ 蘇格拉底式提問：{topic}\n"]
        for i, (layer_name, question) in enumerate(layers, 1):
            lines.append(f"第{i}層【{layer_name}】")
            lines.append(f"  ❓ {question}\n")
        lines.append(f"── 背景參考 ──")
        for h in hits[:2]:
            lines.append(f"• {h.get('title','')}：{h.get('body','')[:150]}")
        lines.append("\n💡 真正的理解來自不斷追問，而非接受第一個答案。")
        return "\n".join(lines)
    except Exception as e:
        return f"蘇格拉底提問失敗：{e}"


def fetch_analogy_maker(concept: str, audience: str = "一般大眾", count: int = 3, lang: str = "zh-tw") -> str:
    try:
        from ddgs import DDGS
        region = "tw-tzh" if lang == "zh-tw" else "us-en"
        with DDGS() as ddgs:
            hits = list(ddgs.text(f"{concept} 類比 比喻 解釋 說明", region=region, max_results=4))
        context = "\n".join(f"{h.get('title','')}：{h.get('body','')[:200]}" for h in hits[:4])
        lines = [
            f"🎯 類比說明：{concept}",
            f"目標受眾：{audience}\n",
            f"── 類比方案 ──\n",
        ]
        # 從搜尋結果找已有的類比，並提示生成
        existing = []
        for h in hits:
            body = h.get("body","")
            if any(kw in body for kw in ["就像", "好比", "類似", "如同", "比喻", "像是", "like", "similar to"]):
                existing.append(body[:200])
        for i, ex in enumerate(existing[:count], 1):
            lines.append(f"類比 {i}：{ex}\n")
        if len(existing) < count:
            lines.append(f"（以上為搜尋到的現有類比，Claude 會在回應中補充更多針對「{audience}」的類比說明）")
        lines.append(f"\n── 原始資料 ──")
        for h in hits[:2]:
            lines.append(f"• {h.get('title','')}：{h.get('body','')[:150]}")
        return "\n".join(lines)
    except Exception as e:
        return f"類比說明失敗：{e}"


def fetch_narrative_builder(topic: str, key_message: str = "", audience: str = "", lang: str = "zh-tw") -> str:
    try:
        from ddgs import DDGS
        region = "tw-tzh" if lang == "zh-tw" else "us-en"
        with DDGS() as ddgs:
            hits = list(ddgs.text(f"{topic} 問題 現況 挑戰 衝突", region=region, max_results=5))
        context = {
            "問題/現況": [],
            "衝突/張力": [],
            "洞察/轉折": [],
            "結論/行動": [],
        }
        for h in hits[:5]:
            body = h.get("body","")[:200]
            title = h.get("title","")
            context["問題/現況"].append(f"{title}：{body}")
        lines = [
            f"📖 敘事架構：{topic}\n",
        ]
        if key_message:
            lines.append(f"核心訊息：{key_message}")
        if audience:
            lines.append(f"目標受眾：{audience}")
        lines.append("")
        lines += [
            f"══ 第一幕：問題／現況 ══",
            f"（建立共鳴，讓受眾認識到問題的存在）",
        ]
        for item in context["問題/現況"][:2]:
            lines.append(f"• {item[:180]}")
        lines += [
            f"\n══ 第二幕：衝突／張力 ══",
            f"（說明為何現有方案不夠，製造戲劇張力）",
            f"• 現有做法的局限：需進一步搜尋或分析",
            f"\n══ 第三幕：洞察／轉折 ══",
            f"（提出新視角或解法，這是故事的核心）",
            f"• 核心洞察：{key_message or '待Claude結合數據補充'}",
            f"\n══ 第四幕：結論／行動 ══",
            f"（清楚的號召行動或結語）",
            f"• 建議行動：根據以上分析，Claude 將在回應中提出具體建議",
        ]
        return "\n".join(lines)
    except Exception as e:
        return f"敘事架構失敗：{e}"


def fetch_critique_writer(subject: str, type_: str = "觀點", lang: str = "zh-tw") -> str:
    try:
        from ddgs import DDGS
        region = "tw-tzh" if lang == "zh-tw" else "us-en"
        queries = [
            f"{subject} 評論 分析 優點",
            f"{subject} 批評 缺點 問題 盲點",
            f"{subject} 背景 脈絡 假設",
        ]
        sections = {}
        with DDGS() as ddgs:
            for q in queries:
                hits = list(ddgs.text(q, region=region, max_results=3))
                sections[q] = hits
        lines = [f"🔬 批判性評析：{subject}（{type_}）\n"]
        lines.append("── 優點與貢獻 ──")
        for h in sections[queries[0]][:2]:
            lines.append(f"✅ {h.get('title','')}：{h.get('body','')[:180]}")
        lines.append("\n── 盲點與缺失 ──")
        for h in sections[queries[1]][:2]:
            lines.append(f"⚠️ {h.get('title','')}：{h.get('body','')[:180]}")
        lines.append("\n── 隱藏假設與脈絡 ──")
        for h in sections[queries[2]][:2]:
            lines.append(f"🔍 {h.get('title','')}：{h.get('body','')[:180]}")
        lines += [
            "\n── 改進建議 ──",
            "（Claude 將根據以上資料，在回應中提出具體改進方向）",
            "\n💡 好的批判不是否定，而是幫助對象看見自己看不見的角落。",
        ]
        return "\n".join(lines)
    except Exception as e:
        return f"批判性評析失敗：{e}"


def fetch_position_statement(issue: str, stance: str, lang: str = "zh-tw") -> str:
    try:
        from ddgs import DDGS
        region = "tw-tzh" if lang == "zh-tw" else "us-en"
        stance_q = "支持 贊成 優點" if stance == "支持" else ("反對 問題 缺點" if stance == "反對" else "條件 但書 前提")
        queries = [
            f"{issue} {stance_q} 論據 證據",
            f"{issue} 數據 研究 案例",
            f"{issue} 反對方 質疑 反駁",
        ]
        evidence, data, counter = [], [], []
        with DDGS() as ddgs:
            evidence = list(ddgs.text(queries[0], region=region, max_results=4))
            data = list(ddgs.text(queries[1], region=region, max_results=3))
            counter = list(ddgs.text(queries[2], region=region, max_results=3))
        lines = [
            f"📣 立場聲明：{issue}",
            f"立場：{stance}\n",
            f"══ 論點（Claim）══",
            f"對於「{issue}」，我的立場是【{stance}】，理由如下：\n",
            f"── 論據與證據 ──",
        ]
        for h in evidence[:3]:
            lines.append(f"• {h.get('title','')}：{h.get('body','')[:180]}")
        lines.append("\n── 數據支撐 ──")
        for h in data[:2]:
            lines.append(f"• {h.get('title','')}：{h.get('body','')[:180]}")
        lines.append("\n── 預判反駁與回應 ──")
        for h in counter[:2]:
            lines.append(f"反方：{h.get('title','')}：{h.get('body','')[:150]}")
        lines += [
            "\n── 結論 ──",
            f"綜合以上，{stance}「{issue}」的立場是有根據且可辯護的。",
            "（Claude 將在回應中整合以上資料，給出完整系統性論述）",
        ]
        return "\n".join(lines)
    except Exception as e:
        return f"立場聲明失敗：{e}"


# ══════════════════════════════════════════════════
#  通用螢幕控制 helper（SendInput，支援負座標/螢幕2）
# ══════════════════════════════════════════════════

def _get_virtual_desktop():
    """回傳 (vl, vt, vw, vh, is_phys, sx, sy)
    vl/vt/vw/vh = 物理虛擬桌面，sx/sy = 邏輯→物理縮放比"""
    import ctypes as _c, mss as _m
    u32 = _c.windll.user32
    old = u32.SetThreadDpiAwarenessContext(_c.c_void_p(-1))
    vl = u32.GetSystemMetrics(76); vt = u32.GetSystemMetrics(77)
    vw = u32.GetSystemMetrics(78); vh = u32.GetSystemMetrics(79)
    u32.SetThreadDpiAwarenessContext(old)
    with _m.mss() as s:
        ms = s.monitors[1:]
        vl_log = min(m["left"] for m in ms)
        vt_log = min(m["top"] for m in ms)
        vw_log = max(m["left"]+m["width"] for m in ms) - vl_log
        vh_log = max(m["top"]+m["height"] for m in ms) - vt_log
    is_phys = vw > 4000
    sx = vw / vw_log if is_phys and vw_log > 0 else 1.0
    sy = vh / vh_log if is_phys and vh_log > 0 else 1.0
    return vl, vt, vw, vh, is_phys, sx, sy

def _si_universal(ax_log: float, ay_log: float, click_type: str = "click"):
    """SendInput 點擊，ax_log/ay_log 為 mss 邏輯座標（支援負值螢幕2）"""
    import ctypes as _c, ctypes.wintypes as _w, time as _t
    vl, vt, vw, vh, is_phys, sx, sy = _get_virtual_desktop()
    ax = round(ax_log * sx); ay = round(ay_log * sy)
    nx = int((ax - vl) * 65535 // vw)
    ny = int((ay - vt) * 65535 // vh)
    u32 = _c.windll.user32
    class MI(_c.Structure):
        _fields_ = [('dx',_w.LONG),('dy',_w.LONG),('mouseData',_w.DWORD),
                    ('dwFlags',_w.DWORD),('time',_w.DWORD),('dwExtraInfo',_c.POINTER(_c.c_ulong))]
    class U(_c.Union): _fields_ = [('mi', MI)]
    class INP(_c.Structure):
        _anonymous_ = ('u',); _fields_ = [('type',_w.DWORD),('u',U)]
    def _send(flags, dx=0, dy=0, md=0):
        i = INP(0, U(mi=MI(dx, dy, md, flags, 0, None)))
        u32.SendInput(1, _c.byref(i), _c.sizeof(i))
    _send(0x0001|0x8000|0x4000, nx, ny)
    _t.sleep(0.25)
    if click_type == "double_click":
        _send(0x0002); _t.sleep(0.05); _send(0x0004); _t.sleep(0.05)
        _send(0x0002); _t.sleep(0.05); _send(0x0004)
    elif click_type == "right_click":
        _send(0x0008); _t.sleep(0.05); _send(0x0010)
    else:
        _send(0x0002); _t.sleep(0.05); _send(0x0004)

def _si_scroll(ax_log: float, ay_log: float, amount: int, direction: str = "down"):
    """SendInput 滾輪，在邏輯座標位置向上/下滾動"""
    import ctypes as _c, ctypes.wintypes as _w, time as _t
    vl, vt, vw, vh, is_phys, sx, sy = _get_virtual_desktop()
    ax = round(ax_log * sx); ay = round(ay_log * sy)
    nx = int((ax - vl) * 65535 // vw)
    ny = int((ay - vt) * 65535 // vh)
    u32 = _c.windll.user32
    class MI(_c.Structure):
        _fields_ = [('dx',_w.LONG),('dy',_w.LONG),('mouseData',_w.DWORD),
                    ('dwFlags',_w.DWORD),('time',_w.DWORD),('dwExtraInfo',_c.POINTER(_c.c_ulong))]
    class U(_c.Union): _fields_ = [('mi', MI)]
    class INP(_c.Structure):
        _anonymous_ = ('u',); _fields_ = [('type',_w.DWORD),('u',U)]
    def _send(flags, dx=0, dy=0, md=0):
        i = INP(0, U(mi=MI(dx, dy, _c.c_ulong(md & 0xFFFFFFFF).value, flags, 0, None)))
        u32.SendInput(1, _c.byref(i), _c.sizeof(i))
    _send(0x0001|0x8000|0x4000, nx, ny)
    _t.sleep(0.15)
    delta = 120 * amount * (1 if direction == "up" else -1)
    for _ in range(max(1, amount // 3)):
        _send(0x0800, md=120 * 3 * (1 if direction == "up" else -1))
        _t.sleep(0.08)


def _cap_monitor_logical(monitor: int):
    """截圖回傳 (PIL.Image, mon_left_log, mon_top_log) 使用 mss 邏輯座標"""
    import mss
    from PIL import Image as _PI
    with mss.mss() as s:
        mon = s.monitors[monitor] if monitor < len(s.monitors) else s.monitors[1]
        shot = s.grab(mon)
        img = _PI.frombytes("RGB", shot.size, shot.bgra, "raw", "BGRX")
    return img, mon["left"], mon["top"]


# ── YOLO UI 偵測器 ────────────────────────────────────────────────────
_yolo_model = None

def _yolo_detect(img, conf=0.4):
    """用 YOLO 偵測螢幕上的物件，回傳 [(label, x_center, y_center, w, h, confidence), ...]
    載入 yolo_ui.pt（通用模型或自訂 UI 模型）
    """
    global _yolo_model
    try:
        if _yolo_model is None:
            from pathlib import Path
            model_path = Path(__file__).parent / "yolo_ui.pt"
            if not model_path.exists():
                return []
            from ultralytics import YOLO
            _yolo_model = YOLO(str(model_path))
        import numpy as np
        from PIL import Image as _PI
        if isinstance(img, _PI.Image):
            img_arr = np.array(img)
        else:
            img_arr = img
        results = _yolo_model(img_arr, conf=conf, verbose=False)
        detections = []
        for r in results:
            for box in r.boxes:
                cls_id = int(box.cls[0])
                label = _yolo_model.names[cls_id]
                x1, y1, x2, y2 = box.xyxy[0].tolist()
                cx, cy = int((x1 + x2) / 2), int((y1 + y2) / 2)
                w, h = int(x2 - x1), int(y2 - y1)
                detections.append((label, cx, cy, w, h, float(box.conf[0])))
        return detections
    except Exception as e:
        return []


# ── 螢幕變化事件流 ───────────────────────────────────────────────────
import threading

class ScreenEventStream:
    """背景持續監控螢幕，偵測變化事件"""

    def __init__(self, monitor=1, fps=2):
        self._monitor = monitor
        self._interval = 1.0 / fps
        self._running = False
        self._thread = None
        self._callbacks = {
            "stable": [],      # 畫面穩定（載入完成）
            "change": [],      # 畫面大幅變化（彈窗/頁面切換）
            "motion": [],      # 持續有動態（影片播放中）
        }
        self._last_frame = None
        self._stable_count = 0

    def on(self, event: str, callback):
        """註冊事件回調：event = 'stable' | 'change' | 'motion'"""
        if event in self._callbacks:
            self._callbacks[event].append(callback)

    def start(self):
        if self._running:
            return
        self._running = True
        self._thread = threading.Thread(target=self._loop, daemon=True)
        self._thread.start()

    def stop(self):
        self._running = False
        if self._thread:
            self._thread.join(timeout=3)

    @property
    def is_running(self):
        return self._running

    def _loop(self):
        import numpy as np, time
        while self._running:
            try:
                img, _, _ = _cap_monitor_logical(self._monitor)
                frame = np.array(img)

                if self._last_frame is not None and frame.shape == self._last_frame.shape:
                    diff = np.mean(np.abs(frame.astype(float) - self._last_frame.astype(float)))
                    diff_pct = diff / 255 * 100

                    if diff_pct < 0.3:
                        self._stable_count += 1
                        if self._stable_count == 3:  # 連續 3 幀穩定
                            for cb in self._callbacks["stable"]:
                                try: cb(diff_pct)
                                except Exception: pass
                    elif diff_pct > 5.0:
                        self._stable_count = 0
                        for cb in self._callbacks["change"]:
                            try: cb(diff_pct)
                            except Exception: pass
                    else:
                        self._stable_count = 0
                        for cb in self._callbacks["motion"]:
                            try: cb(diff_pct)
                            except Exception: pass

                self._last_frame = frame
            except Exception:
                pass
            time.sleep(self._interval)

    def wait_for_stable(self, timeout=15.0):
        """阻塞等待畫面穩定，回傳 True=穩定 / False=超時"""
        import time
        event = threading.Event()
        def _on_stable(diff):
            event.set()
        self.on("stable", _on_stable)
        was_running = self._running
        if not was_running:
            self.start()
        result = event.wait(timeout=timeout)
        if not was_running:
            self.stop()
        # 清掉這個一次性回調
        try: self._callbacks["stable"].remove(_on_stable)
        except ValueError: pass
        return result

    def wait_for_change(self, timeout=30.0):
        """阻塞等待畫面出現大幅變化"""
        import time
        event = threading.Event()
        def _on_change(diff):
            event.set()
        self.on("change", _on_change)
        was_running = self._running
        if not was_running:
            self.start()
        result = event.wait(timeout=timeout)
        if not was_running:
            self.stop()
        try: self._callbacks["change"].remove(_on_change)
        except ValueError: pass
        return result


# ── 操作狀態機 ───────────────────────────────────────────────────────
class OperationStateMachine:
    """多步驟操作的狀態機，每一步有成功/失敗判斷和重試"""

    def __init__(self, name="operation"):
        self.name = name
        self.steps = []
        self.current = 0
        self.results = []

    def add_step(self, name: str, action, verify=None, max_retries=3, on_fail="retry"):
        """加入一步
        action: callable，執行動作，回傳 any
        verify: callable(result) → bool，驗證成功與否（None=不驗證）
        on_fail: "retry" | "skip" | "abort"
        """
        self.steps.append({
            "name": name,
            "action": action,
            "verify": verify,
            "max_retries": max_retries,
            "on_fail": on_fail,
        })

    def run(self) -> dict:
        """執行所有步驟，回傳 {"ok": bool, "results": [...], "failed_step": str|None}"""
        self.current = 0
        self.results = []
        for i, step in enumerate(self.steps):
            self.current = i
            success = False
            result = None
            for attempt in range(step["max_retries"]):
                try:
                    result = step["action"]()
                    if step["verify"] is None:
                        success = True
                        break
                    elif step["verify"](result):
                        success = True
                        break
                except Exception as e:
                    result = str(e)
                import time; time.sleep(0.5)

            self.results.append({
                "step": step["name"],
                "success": success,
                "result": result,
            })

            if not success:
                if step["on_fail"] == "abort":
                    return {"ok": False, "results": self.results, "failed_step": step["name"]}
                elif step["on_fail"] == "skip":
                    continue
                else:  # retry exhausted
                    return {"ok": False, "results": self.results, "failed_step": step["name"]}

        return {"ok": True, "results": self.results, "failed_step": None}


# 預定義常用操作流程
def youtube_play_flow(keyword: str, monitor: int = 2) -> dict:
    """YouTube 搜尋並播放的完整狀態機流程"""
    import webbrowser, time as _t_yt

    sm = OperationStateMachine("youtube_play")

    # Step 1: 開 YouTube 搜尋頁
    def _open():
        webbrowser.open(f"https://www.youtube.com/results?search_query={keyword.replace(' ', '+')}")
        return True
    sm.add_step("open_search", _open, max_retries=2, on_fail="abort")

    # Step 2: 等頁面載入
    def _wait_load():
        _t_yt.sleep(1)  # 給瀏覽器反應時間
        return _wait_screen_stable(monitor, threshold=0.5, timeout=10.0)
    def _verify_load(result):
        return result  # True = 穩定了
    sm.add_step("wait_load", _wait_load, _verify_load, max_retries=2, on_fail="abort")

    # Step 3: 滾過廣告
    def _scroll():
        import pyautogui, win32gui, win32con, ctypes
        wins = []
        def _fb(h, _):
            if win32gui.IsWindowVisible(h):
                t = win32gui.GetWindowText(h).lower()
                if 'youtube' in t or 'chrome' in t:
                    wins.append(h)
            return True
        win32gui.EnumWindows(_fb, None)
        if wins:
            ctypes.windll.user32.SetForegroundWindow(wins[0])
            _t_yt.sleep(0.5)
        pyautogui.scroll(-3)
        _t_yt.sleep(1)
        return True
    sm.add_step("scroll_past_ads", _scroll, max_retries=1, on_fail="skip")

    # Step 4: vision_locate 點擊影片
    def _click_video():
        return fetch_vision_locate(
            f"YouTube搜尋結果中{keyword}的官方MV或歌曲影片縮圖（有OFFICIAL標誌的優先，跳過贊助商廣告）",
            monitor, "click"
        )
    def _verify_click(result):
        return "找到" in str(result) and "找不到" not in str(result)
    sm.add_step("click_video", _click_video, _verify_click, max_retries=3, on_fail="abort")

    # Step 5: 等影片開始播放
    def _wait_play():
        _t_yt.sleep(1)
        return _wait_screen_stable(monitor, threshold=0.3, timeout=5.0)
    sm.add_step("wait_play", _wait_play, max_retries=1, on_fail="skip")

    return sm.run()


# ── 智慧等待：偵測螢幕停止變化 ─────────────────────────────────────────
def _wait_screen_stable(monitor: int = 1, threshold: float = 0.5, timeout: float = 15.0, interval: float = 0.5):
    """連續截圖比對，直到畫面穩定（差異低於 threshold%）或超時"""
    import time, numpy as np
    img_a, _, _ = _cap_monitor_logical(monitor)
    arr_a = np.array(img_a)
    start = time.time()
    while time.time() - start < timeout:
        time.sleep(interval)
        img_b, _, _ = _cap_monitor_logical(monitor)
        arr_b = np.array(img_b)
        if arr_a.shape == arr_b.shape:
            diff = np.mean(np.abs(arr_a.astype(float) - arr_b.astype(float)))
            diff_pct = diff / 255 * 100
            if diff_pct < threshold:
                return True  # 畫面穩定
        arr_a = arr_b
    return False  # 超時


# ── 動作驗證迴圈：點擊後確認畫面有變化 ────────────────────────────────
def _action_with_verify(action_fn, monitor: int = 1, max_retries: int = 3, wait: float = 0.8):
    """執行動作 → 比對前後截圖 → 沒變化就重試
    action_fn: 無參數的 callable，執行實際動作（如點擊）
    回傳 True 如果畫面有變化，False 如果重試後仍無變化
    """
    import time, numpy as np
    for attempt in range(max_retries):
        img_before, _, _ = _cap_monitor_logical(monitor)
        arr_before = np.array(img_before)
        action_fn()
        time.sleep(wait)
        img_after, _, _ = _cap_monitor_logical(monitor)
        arr_after = np.array(img_after)
        if arr_before.shape == arr_after.shape:
            diff = np.mean(np.abs(arr_before.astype(float) - arr_after.astype(float)))
            diff_pct = diff / 255 * 100
            if diff_pct > 0.5:  # 畫面有變化 = 動作成功
                return True
        # 沒變化，重試
        time.sleep(0.3)
    return False


# ── UIA 元素偵測：用 Windows UI Automation 找元素 ─────────────────────
def _uia_find_element(description: str, window_title: str = None):
    """用 pywinauto UIA 在前景視窗中搜尋元素，回傳 (center_x, center_y) 或 (None, None)
    比 vision_locate 快 10 倍，不需截圖和 API
    """
    try:
        from pywinauto import Desktop
        import re

        desktop = Desktop(backend="uia")

        # 找目標視窗
        if window_title:
            wins = desktop.windows(title_re=f".*{re.escape(window_title)}.*", visible_only=True)
        else:
            # 用前景視窗
            import win32gui
            hwnd = win32gui.GetForegroundWindow()
            wins = [desktop.window(handle=hwnd)]

        if not wins:
            return None, None

        win = wins[0]
        desc_lower = description.lower()

        # 搜尋策略：先精確匹配，再模糊匹配
        best = None
        best_score = 0

        try:
            # 取得所有子元素（限深度避免太慢）
            children = win.descendants(depth=8)
        except Exception:
            return None, None

        for child in children:
            try:
                name = (child.window_text() or "").strip()
                ctrl_type = (child.friendly_class_name() or "").lower()
                if not name:
                    continue

                name_lower = name.lower()
                score = 0

                # 精確包含
                if desc_lower in name_lower:
                    score = 100
                # 關鍵字匹配
                else:
                    keywords = [w for w in desc_lower.split() if len(w) > 1]
                    if keywords:
                        matched = sum(1 for kw in keywords if kw in name_lower)
                        score = matched / len(keywords) * 80

                # 可點擊元素加分
                if ctrl_type in ("button", "hyperlink", "listitem", "menuitem", "treeitem", "tabitem"):
                    score += 10

                if score > best_score:
                    rect = child.rectangle()
                    cx = (rect.left + rect.right) // 2
                    cy = (rect.top + rect.bottom) // 2
                    # 確保座標在合理範圍
                    if -4000 < cx < 8000 and -2000 < cy < 4000:
                        best = (cx, cy)
                        best_score = score
            except Exception:
                continue

        if best and best_score >= 50:
            return best
        return None, None
    except Exception:
        return None, None


# ── YOLO 訓練資料自動收集 ──────────────────────────────────────────
def _save_training_sample(img, description: str, rx: int, ry: int):
    """每次 vision_locate 成功時自動保存截圖+標記，累積 YOLO 訓練資料"""
    try:
        from pathlib import Path
        import json, time
        data_dir = Path(__file__).parent / "yolo_training_data"
        data_dir.mkdir(exist_ok=True)
        ts = int(time.time() * 1000)
        img.save(data_dir / f"{ts}.jpg", quality=90)
        # YOLO 格式：class x_center y_center width height（歸一化）
        # 暫時用固定 50x50 框，之後可以調整
        w, h = img.width, img.height
        box_w, box_h = 50 / w, 50 / h
        cx, cy = rx / w, ry / h
        with open(data_dir / f"{ts}.txt", "w") as f:
            f.write(f"0 {cx:.6f} {cy:.6f} {box_w:.6f} {box_h:.6f}\n")
        with open(data_dir / f"{ts}.json", "w", encoding="utf-8") as f:
            json.dump({"description": description, "x": rx, "y": ry, "img_w": w, "img_h": h}, f, ensure_ascii=False)
    except Exception:
        pass


def _vision_find(img, description: str):
    """截圖 → OCR輔助 + Claude Vision → 回傳 (rx, ry) resized圖像像素，找不到回傳 (None, None)"""
    import anthropic, base64, io, json, re
    from PIL import Image as _PI
    ow, oh = img.width, img.height
    # 2048px / quality 92（于晏哥指定的辨識品質標準）
    if img.width > 2048:
        r = 2048 / img.width
        img = img.resize((2048, int(img.height * r)), _PI.LANCZOS)
    scale = ow / img.width
    buf = io.BytesIO(); img.save(buf, format="JPEG", quality=92)
    # 超過 4MB 則降品質重試
    if buf.tell() > 4 * 1024 * 1024:
        buf = io.BytesIO(); img.save(buf, format="JPEG", quality=80)
    b64 = base64.standard_b64encode(buf.getvalue()).decode()

    # OCR 輔助：提取螢幕上的文字供 Claude 參考
    ocr_hint = ""
    try:
        import pytesseract
        ocr_data = pytesseract.image_to_data(img, lang="chi_tra+eng", output_type=pytesseract.Output.DICT)
        ocr_items = []
        for i, txt in enumerate(ocr_data["text"]):
            if txt.strip() and ocr_data["conf"][i] > 40:
                cx = ocr_data["left"][i] + ocr_data["width"][i] // 2
                cy = ocr_data["top"][i] + ocr_data["height"][i] // 2
                ocr_items.append(f"「{txt.strip()}」at({cx},{cy})")
        if ocr_items:
            ocr_hint = "\n\nOCR偵測到的文字及位置：" + "; ".join(ocr_items[:30])
    except Exception:
        pass

    prompt = (
        f"你是專業的螢幕分析師。這是一張電腦螢幕截圖（{img.width}x{img.height}px）。\n"
        f"請找到「{description}」的中心座標。\n"
        f"注意：\n"
        f"- 座標是相對於圖片左上角(0,0)的像素位置\n"
        f"- 仔細區分廣告和真正的內容（廣告通常有「贊助商廣告」「Ad」「探索」字樣）\n"
        f"- 如果有多個匹配，選最相關的那個\n"
        f"- 僅回傳JSON: {{\"x\":整數, \"y\":整數, \"ok\":true/false}}"
        f"{ocr_hint}"
    )

    resp = anthropic.Anthropic().messages.create(
        model="claude-sonnet-4-6", max_tokens=300,
        messages=[{"role": "user", "content": [
            {"type": "image", "source": {"type": "base64", "media_type": "image/jpeg", "data": b64}},
            {"type": "text", "text": prompt}
        ]}]
    )
    m = re.search(r'\{.*?\}', resp.content[0].text, re.DOTALL)
    if not m: return None, None
    d = json.loads(m.group())
    if not d.get("ok", True): return None, None
    rx, ry = int(d["x"] * scale), int(d["y"] * scale)
    # 自動收集訓練資料
    _save_training_sample(img, description, d["x"], d["y"])
    return rx, ry


def fetch_ocr_click(target_text: str, monitor: int = 1, click_type: str = "click", region: list = None) -> str:
    try:
        import pyautogui
        import pytesseract
        from PIL import Image
        import numpy as np

        # 截圖
        mon_map = {1: 0, 2: 1, 3: 2}
        try:
            import dxcam
            cam = dxcam.create(device_idx=mon_map.get(monitor, 0))
            frame = cam.grab()
            cam.release()
            if frame is None:
                raise Exception("dxcam grab failed")
            img = Image.fromarray(frame)
        except Exception:
            import mss
            with mss.mss() as sct:
                monitors = sct.monitors
                mon = monitors[monitor] if monitor < len(monitors) else monitors[1]
                if region:
                    mon = {"left": mon["left"] + region[0], "top": mon["top"] + region[1],
                           "width": region[2], "height": region[3]}
                shot = sct.grab(mon)
                img = Image.frombytes("RGB", shot.size, shot.bgra, "raw", "BGRX")

        # OCR
        data = pytesseract.image_to_data(img, lang="chi_tra+eng", output_type=pytesseract.Output.DICT)
        found = []
        for i, text in enumerate(data["text"]):
            if target_text.lower() in text.lower() and data["conf"][i] > 30:
                x = data["left"][i] + data["width"][i] // 2
                y = data["top"][i] + data["height"][i] // 2
                found.append((x, y, data["conf"][i], text))

        if not found:
            return f"OCR找不到文字「{target_text}」，請確認文字正確或改用 vision_locate"

        # 取信心最高的結果，換算為螢幕絕對座標
        best = max(found, key=lambda f: f[2])
        cx, cy = best[0], best[1]

        # 加上螢幕偏移
        try:
            import mss
            with mss.mss() as sct:
                mon = sct.monitors[monitor] if monitor < len(sct.monitors) else sct.monitors[1]
                abs_x = mon["left"] + cx
                abs_y = mon["top"] + cy
        except Exception:
            abs_x, abs_y = cx, cy

        # 執行點擊（SendInput，支援螢幕2負座標）
        _si_universal(abs_x, abs_y, click_type)
        return f"✅ OCR找到「{best[3]}」（信心{best[2]}%），已在 ({abs_x}, {abs_y}) 執行 {click_type}"
    except Exception as e:
        return f"OCR點擊失敗：{e}"


def fetch_vision_locate(description: str, monitor: int = 1, action: str = "click", region: list = None) -> str:
    try:
        # 嘗試把瀏覽器切到前景並最大化（如果描述中提到影片、YouTube等）
        _desc_lower = description.lower()
        _browser_kw = any(k in _desc_lower for k in ["影片", "video", "youtube", "縮圖", "thumbnail", "搜尋結果", "第一", "chrome", "網頁"])
        if _browser_kw:
            try:
                import win32gui, win32con
                def _find_chrome(h, results):
                    if win32gui.IsWindowVisible(h):
                        t = win32gui.GetWindowText(h).lower()
                        if "youtube" in t or "chrome" in t or "edge" in t or "firefox" in t:
                            results.append((h, t))
                    return True
                _browser_wins = []
                win32gui.EnumWindows(_find_chrome, _browser_wins)
                _yt_wins = [h for h, t in _browser_wins if "youtube" in t]
                _target = _yt_wins[0] if _yt_wins else (_browser_wins[0][0] if _browser_wins else None)
                if _target:
                    win32gui.ShowWindow(_target, win32con.SW_RESTORE)
                    win32gui.ShowWindow(_target, win32con.SW_MAXIMIZE)
                    import ctypes
                    ctypes.windll.user32.SetForegroundWindow(_target)
                    import time; time.sleep(0.5)
            except Exception:
                pass

        # ── 策略1：UIA 快速搜尋（不需截圖和 API，最快） ──
        if not region and action != "locate_only":
            uia_x, uia_y = _uia_find_element(description)
            if uia_x is not None:
                if action != "locate_only":
                    # 用動作驗證迴圈點擊
                    _det_mon = monitor if not _browser_kw else (2 if monitor == 1 else monitor)
                    success = _action_with_verify(
                        lambda: _si_universal(uia_x, uia_y, action),
                        monitor=_det_mon
                    )
                    v = "已驗證" if success else "未驗證"
                    return f"✅ UIA找到「{description}」，已在 ({uia_x}, {uia_y}) 執行 {action}（{v}）"
                return f"✅ UIA找到「{description}」，位置 ({uia_x}, {uia_y})"

        # ── 策略2：YOLO 快速偵測（本地，0.05秒） ──
        if not region:
            try:
                img_yolo, mon_left_y, mon_top_y = _cap_monitor_logical(monitor)
                detections = _yolo_detect(img_yolo, conf=0.3)
                if detections:
                    # 用描述關鍵字匹配 YOLO 偵測到的物件
                    _desc_words = [w.lower() for w in description.split() if len(w) > 1]
                    best_det = None
                    best_score = 0
                    for label, cx, cy, w, h, conf in detections:
                        score = sum(1 for kw in _desc_words if kw in label.lower()) * 50 + conf * 30
                        if score > best_score:
                            best_det = (label, cx, cy, conf)
                            best_score = score
                    if best_det and best_score > 30:
                        _yl, _yx, _yy, _yc = best_det
                        abs_x = mon_left_y + _yx
                        abs_y = mon_top_y + _yy
                        if action == "locate_only":
                            return f"✅ YOLO找到「{_yl}」({_yc:.0%})，位置 ({abs_x}, {abs_y})"
                        success = _action_with_verify(
                            lambda: _si_universal(abs_x, abs_y, action),
                            monitor=monitor
                        )
                        v = "已驗證" if success else "未驗證"
                        return f"✅ YOLO找到「{_yl}」({_yc:.0%})，已在 ({abs_x}, {abs_y}) 執行 {action}（{v}）"
            except Exception:
                pass

        # ── 策略3：智慧等待 + 視覺辨識（截圖 + Sonnet） ──
        # 等螢幕穩定再截圖
        _wait_screen_stable(monitor, threshold=0.5, timeout=3.0, interval=0.3)

        img, mon_left, mon_top = _cap_monitor_logical(monitor)

        if region:
            from PIL import Image as _PI
            img = img.crop((region[0], region[1], region[0]+region[2], region[1]+region[3]))
            mon_left += region[0]; mon_top += region[1]

        rx, ry = _vision_find(img, description)
        if rx is None:
            return f"視覺辨識找不到「{description}」"

        abs_x = mon_left + rx
        abs_y = mon_top + ry

        if action == "locate_only":
            return f"✅ 找到「{description}」，位置 ({abs_x}, {abs_y})"

        # 用動作驗證迴圈點擊
        success = _action_with_verify(
            lambda: _si_universal(abs_x, abs_y, action),
            monitor=monitor
        )
        v = "已驗證" if success else "未驗證"
        return f"✅ 視覺找到「{description}」，已在 ({abs_x}, {abs_y}) 執行 {action}（{v}）"
    except Exception as e:
        return f"視覺定位失敗：{e}"


def fetch_screen_workflow(steps: list) -> str:
    try:
        import time
        results = []
        for i, step in enumerate(steps):
            action = step.get("action", "")
            target = step.get("target", "")
            value = step.get("value", "")
            monitor = step.get("monitor", 1)
            try:
                if action == "screenshot":
                    results.append(f"步驟{i+1} screenshot：已截圖")
                elif action == "ocr_click":
                    r = fetch_ocr_click(target, monitor)
                    results.append(f"步驟{i+1} ocr_click [{target}]：{r}")
                elif action == "vision_click":
                    r = fetch_vision_locate(target, monitor, "click")
                    results.append(f"步驟{i+1} vision_click [{target}]：{r}")
                elif action == "type":
                    import pyautogui
                    pyautogui.write(value, interval=0.05)
                    results.append(f"步驟{i+1} type：已輸入「{value[:30]}」")
                elif action == "press":
                    import pyautogui
                    pyautogui.press(value)
                    results.append(f"步驟{i+1} press：已按 {value}")
                elif action == "wait":
                    secs = float(value) if value else 1.0
                    time.sleep(secs)
                    results.append(f"步驟{i+1} wait：等待 {secs}s")
                elif action == "open_app":
                    import subprocess
                    subprocess.Popen(target, shell=True)
                    time.sleep(1.5)
                    results.append(f"步驟{i+1} open_app：已開啟 {target}")
                elif action == "hotkey":
                    import pyautogui
                    keys = [k.strip() for k in value.split("+")]
                    pyautogui.hotkey(*keys)
                    results.append(f"步驟{i+1} hotkey：{value}")
                else:
                    results.append(f"步驟{i+1} 未知動作：{action}")
            except Exception as e:
                results.append(f"步驟{i+1} 失敗：{e}")
                break
        return "📋 工作流執行結果：\n" + "\n".join(results)
    except Exception as e:
        return f"螢幕工作流失敗：{e}"


def fetch_app_navigator(app: str, task: str, input_text: str = "", monitor: int = 1, contact_name: str = "") -> str:
    try:
        import time
        import pyautogui
        import re
        app_lower = app.lower()
        results = []

        # 取得螢幕絕對座標偏移
        try:
            import mss
            with mss.mss() as sct:
                mons = sct.monitors
                mon_info = mons[monitor] if monitor < len(mons) else mons[1]
                mon_left = mon_info["left"]
                mon_top = mon_info["top"]
                mon_cx = mon_left + mon_info["width"] // 2
                mon_cy = mon_top + mon_info["height"] // 2
        except Exception:
            mon_left, mon_top, mon_cx, mon_cy = 0, 0, 960, 540

        # ══ Telegram 純螢幕控制（最優先，跳過所有 win32gui 視窗管理）══
        if "telegram" in app_lower:
            import pyperclip as _pc
            import ctypes as _ct, ctypes.wintypes as _wt
            import anthropic as _ant, base64 as _b64, io as _io2
            import ctypes as _ct2, win32ui as _w32u, win32con as _w32c2, win32gui as _w32g2
            import mss as _mss2
            from PIL import Image as _PI

            # 提取聯絡人名稱（優先用 contact_name 直接參數，再引號，再 regex）
            if contact_name:
                name = contact_name.strip()
            else:
                name_match = re.search(r'[「"](.*?)[」"]', task)
                if name_match:
                    name = name_match.group(1)
                else:
                    # 優先：「找/搜尋 XXX」格式，停於並/和/的/給/，/說/聊/打
                    m2 = re.search(r'(?:找到?|搜尋)(.+?)(?:並|和|的|給|，|,|說|聊|打|$)', task)
                    if m2:
                        name = m2.group(1).strip()
                    else:
                        # 次選：「跟 XXX 聊/說/打」格式，提取跟和動詞之間的名稱
                        m3 = re.search(r'跟(.+?)(?:聊|說|打|和|並|的|給|，|,|$)', task)
                        if m3:
                            name = m3.group(1).strip()
                        else:
                            name = re.sub(r'(螢幕\d|從Telegram|從telegram|Telegram|telegram|打開|找到|找|搜尋|對話|訊息|聊天|跟他說.*|跟.*?說.*|和|給|傳)', '', task).strip()

            # 名稱為空就直接報錯，不貼錯內容進搜尋框
            if not name:
                return "❌ 無法提取聯絡人名稱，請用引號標明，例如「去找「巴斯」聊天」或加 contact_name 參數"

            # ── 虛擬桌面參數（強制物理 DPI context，SendInput 用物理座標）──
            _u32 = _ct.windll.user32
            _old_ctx = _u32.SetThreadDpiAwarenessContext(_ct.c_void_p(-1))
            _vl = _u32.GetSystemMetrics(76); _vt = _u32.GetSystemMetrics(77)
            _vw = _u32.GetSystemMetrics(78); _vh = _u32.GetSystemMetrics(79)
            _u32.SetThreadDpiAwarenessContext(_old_ctx)
            # 取得螢幕邏輯尺寸，計算 DPI scale
            with _mss2.mss() as _smss:
                _mon_log = _smss.monitors[monitor]
                _log_left = _mon_log["left"]; _log_top = _mon_log["top"]
                _log_w = _mon_log["width"];   _log_h = _mon_log["height"]

            def _si_click(ax, ay):
                # ax, ay 是物理絕對座標（與 GetSystemMetrics DPI-aware 一致）
                class _MI(_ct.Structure):
                    _fields_ = [('dx',_wt.LONG),('dy',_wt.LONG),('mouseData',_wt.DWORD),
                                 ('dwFlags',_wt.DWORD),('time',_wt.DWORD),
                                 ('dwExtraInfo',_ct.POINTER(_ct.c_ulong))]
                class _U(_ct.Union):
                    _fields_ = [('mi',_MI)]
                class _INP(_ct.Structure):
                    _anonymous_ = ('u',); _fields_ = [('type',_wt.DWORD),('u',_U)]
                def _send(flags, dx=0, dy=0):
                    i = _INP(0,_U(mi=_MI(dx,dy,0,flags,0,None)))
                    _u32.SendInput(1,_ct.byref(i),_ct.sizeof(i))
                nx = int((ax - _vl) * 65535 // _vw)
                ny = int((ay - _vt) * 65535 // _vh)
                _send(0x0001|0x8000|0x4000, nx, ny)  # MOVE|ABSOLUTE|VIRTUALDESK
                time.sleep(0.35)
                _send(0x0002); time.sleep(0.08); _send(0x0004)  # DOWN + UP

            # ── GDI BitBlt 截圖（螢幕2用，其他用 dxcam）──
            def _cap(mon_num):
                with _mss2.mss() as s:
                    m = s.monitors[mon_num]
                    ml, mt, mw, mh = m["left"], m["top"], m["width"], m["height"]
                if mon_num == 2:
                    u32 = _ct2.windll.user32
                    old = u32.SetThreadDpiAwarenessContext(_ct2.c_void_p(-1))
                    try:
                        hd = _w32g2.GetDesktopWindow()
                        hdc = _w32g2.GetWindowDC(hd)
                        mdc = _w32u.CreateDCFromHandle(hdc)
                        sdc = mdc.CreateCompatibleDC()
                        bm = _w32u.CreateBitmap()
                        bm.CreateCompatibleBitmap(mdc, mw, mh)
                        sdc.SelectObject(bm)
                        sdc.BitBlt((0,0),(mw,mh),mdc,(ml,mt),_w32c2.SRCCOPY)
                        inf = bm.GetInfo(); bits = bm.GetBitmapBits(True)
                        img = _PI.frombuffer("RGB",(inf["bmWidth"],inf["bmHeight"]),bits,"raw","BGRX",0,1)
                        _w32g2.DeleteObject(bm.GetHandle())
                        sdc.DeleteDC(); mdc.DeleteDC(); _w32g2.ReleaseDC(hd, hdc)
                    finally:
                        u32.SetThreadDpiAwarenessContext(old)
                else:
                    import dxcam as _dx2
                    c = _dx2.create(output_idx={1:0,3:1}.get(mon_num,0))
                    img = _PI.fromarray(c.grab()); del c
                # 回傳 img + 邏輯偏移 + 邏輯寬（供計算 DPI scale）
                return img, ml, mt, mw

            # 判斷 GetSystemMetrics 是否物理座標
            _is_phys_gm = _vw > 4000

            # ── 截圖 → Claude Vision → 回傳與 GetSystemMetrics 同空間的絕對座標 ──
            def _see(prompt, mon_num):
                img, off_x_log, off_y_log, log_w = _cap(mon_num)
                # DPI scale：GDI 物理圖像寬 / mss 邏輯寬
                dpi = img.width / log_w  # 例：2400/1920=1.25
                ow = img.width
                if img.width > 1280:
                    r = 1280 / img.width
                    img = img.resize((1280, int(img.height*r)), _PI.LANCZOS)
                scale = ow / img.width  # 物理像素/resized像素
                buf = _io2.BytesIO()
                img.save(buf, format="JPEG", quality=85)
                b = _b64.standard_b64encode(buf.getvalue()).decode()
                resp = _ant.Anthropic().messages.create(
                    model="claude-haiku-4-5-20251001", max_tokens=80,
                    messages=[{"role":"user","content":[
                        {"type":"image","source":{"type":"base64","media_type":"image/jpeg","data":b}},
                        {"type":"text","text":f"截圖({img.width}x{img.height}px)。找「{prompt}」中心座標。僅回傳JSON:{{\"x\":整數,\"y\":整數,\"ok\":true/false}}"}
                    ]}]
                )
                import json as _j, re as _r2
                m2 = _r2.search(r'\{.*?\}', resp.content[0].text, _r2.DOTALL)
                if not m2: return None, None
                d = _j.loads(m2.group())
                if not d.get("ok", True): return None, None
                if _is_phys_gm:
                    # GetSystemMetrics=物理 → 回傳物理絕對
                    return round(off_x_log*dpi) + int(d["x"]*scale), \
                           round(off_y_log*dpi) + int(d["y"]*scale)
                else:
                    # GetSystemMetrics=邏輯 → pixel/dpi 轉邏輯再加邏輯偏移
                    return off_x_log + int(d["x"]*scale/dpi), \
                           off_y_log + int(d["y"]*scale/dpi)

            # ══ Step 0：先用 Vision 找 Telegram 視窗邊界（不管移到哪都能定位）══
            # Telegram 視覺特徵：左側深色聊天列表側欄、頂部≡選單+搜尋圖示、藍色紙飛機Logo
            def _find_tg_window(mon_num):
                """回傳 Telegram 視窗在螢幕上的絕對座標 (win_x, win_y, win_w, win_h)，失敗回傳 None"""
                img, off_x_log, off_y_log, log_w = _cap(mon_num)
                dpi = img.width / log_w
                ow = img.width
                if img.width > 1280:
                    r = 1280 / img.width
                    img = img.resize((1280, int(img.height*r)), _PI.LANCZOS)
                scale = ow / img.width
                buf = _io2.BytesIO(); img.save(buf, format="JPEG", quality=85)
                b = _b64.standard_b64encode(buf.getvalue()).decode()
                resp = _ant.Anthropic().messages.create(
                    model="claude-haiku-4-5-20251001", max_tokens=120,
                    messages=[{"role":"user","content":[
                        {"type":"image","source":{"type":"base64","media_type":"image/jpeg","data":b}},
                        {"type":"text","text":
                            f"截圖({img.width}x{img.height}px)。"
                            "找 Telegram 桌面應用程式視窗。"
                            "識別特徵：左側有深色聊天列表側欄（含聯絡人名稱和頭像）、頂部有≡漢堡選單圖示和搜尋放大鏡、藍色紙飛機Logo、右側有聊天對話區域。"
                            "回傳視窗邊界 JSON（像素座標）：{\"x\":左上角x, \"y\":左上角y, \"w\":視窗寬度, \"h\":視窗高度, \"ok\":true/false}"}
                    ]}]
                )
                m2 = re.search(r'\{.*?\}', resp.content[0].text, re.DOTALL)
                if not m2: return None
                import json as _j2
                d = _j2.loads(m2.group())
                if not d.get("ok", True): return None
                # 轉換到與 GetSystemMetrics 同空間的絕對座標
                wx = int(d["x"] * scale); wy = int(d["y"] * scale)
                ww = int(d["w"] * scale); wh = int(d["h"] * scale)
                if _is_phys_gm:
                    return round(off_x_log*dpi)+wx, round(off_y_log*dpi)+wy, ww, wh
                else:
                    return off_x_log+int(wx/dpi), off_y_log+int(wy/dpi), int(ww/dpi), int(wh/dpi)

            _tg_win = _find_tg_window(monitor)
            if _tg_win:
                _tw_x, _tw_y, _tw_w, _tw_h = _tg_win
                results.append(f"📐 Telegram視窗({_tw_x},{_tw_y}) {_tw_w}x{_tw_h}")
            else:
                # 找不到視窗邊界時，退回用 monitor 左上角估算
                _tw_x = round(_log_left * 1.25) if _is_phys_gm else _log_left
                _tw_y = round(_log_top * 1.25) if _is_phys_gm else _log_top
                _tw_w = round(_log_w * 1.25) if _is_phys_gm else _log_w
                _tw_h = round(_log_h * 1.25) if _is_phys_gm else _log_h
                results.append(f"⚠️ 找不到視窗邊界，用螢幕偏移估算({_tw_x},{_tw_y})")

            # ── Step 0.5：確保 Telegram 在正常聊天列表狀態 ──
            import pyautogui as _pg
            _pg.press("escape"); time.sleep(0.3)
            _pg.press("escape"); time.sleep(0.3)
            # 檢查是否在異常狀態（設定面板、空白畫面等）
            try:
                import pytesseract as _tess_st
                _state_img, _, _ = _cap(monitor)
                _state_text = _tess_st.image_to_string(_state_img, lang="chi_tra+eng")
                if "請選擇聊天對象" in _state_text or "我的資料" in _state_text or "設定" in _state_text or "建立群組" in _state_text:
                    _pg.press("escape"); time.sleep(0.3)
                    _pg.press("escape"); time.sleep(0.3)
                    results.append("🔄 從異常狀態恢復")
            except Exception:
                pass

            # ① 用鍵盤快捷鍵開搜尋（不用 Vision 找搜尋框，避免點到漢堡選單）
            _si_click(_tw_x + int(_tw_w * 0.13), _tw_y + int(_tw_h * 0.05))
            time.sleep(0.3)
            results.append("🔍 點擊搜尋區域")

            # ② 清空 + 貼上聯絡人名稱
            _pg.hotkey("ctrl","a"); time.sleep(0.1); _pg.press("delete"); time.sleep(0.1)
            _pc.copy(name); _pg.hotkey("ctrl","v"); time.sleep(1.2)
            results.append(f"🔍 搜尋名稱：{name}")

            # ③ 用 Vision 在搜尋結果中找聯絡人
            _rx, _ry = _see(
                f"Telegram搜尋結果列表中名稱為「{name}」的聊天項目（個人對話優先），點名字或頭像位置",
                monitor
            )
            if _rx is not None:
                _si_click(_rx, _ry); time.sleep(0.7)
                results.append(f"✅ Vision找到並點擊「{name}」({_rx},{_ry})")
            else:
                # 備用：搜尋框下方固定偏移點第一個結果
                _first_y = _tw_y + int(_tw_h * 0.11)
                _si_click(_tw_x + int(_tw_w * 0.13), _first_y); time.sleep(0.7)
                results.append(f"⚠️ Vision找不到，點第一個搜尋結果")

            # ④ 有訊息就輸入送出
            if input_text:
                time.sleep(0.5)
                mx, my = _see(
                    "Telegram聊天視窗底部的訊息輸入框（最下方打字區，顯示「輸入訊息」或有emoji圖示和迴紋針圖示，不是頂部搜尋框）",
                    monitor
                )
                if mx is None:
                    mx = _tw_x + int(_tw_w * 0.6); my = _tw_y + int(_tw_h * 0.95)
                    results.append(f"⚠️ Vision找不到訊息框，視窗相對備用({mx},{my})")
                else:
                    _si_click(mx, my); time.sleep(0.4)

                # 偵測編號清單（1. xxx\n2. xxx），拆成多則分開發送
                _lines = re.split(r'\n(?=\d+[\.\、\)])', input_text.strip())
                if len(_lines) > 1:
                    for _li, _line in enumerate(_lines):
                        _msg = re.sub(r'^\d+[\.\、\)]\s*', '', _line).strip()
                        if not _msg:
                            continue
                        _pc.copy(_msg); _pg.hotkey("ctrl","v"); time.sleep(0.2)
                        _pg.press("enter"); time.sleep(0.4)
                        results.append(f"📤 第{_li+1}則：{_msg[:20]}")
                else:
                    _pc.copy(input_text); _pg.hotkey("ctrl","v"); time.sleep(0.3)
                    _pg.press("enter"); time.sleep(0.2); _pg.press("enter")
                    results.append(f"📤 已送：{input_text}")

            # ── 清理：點擊聊天列表第一個聊天回到正常狀態（不用 Esc 避免關掉聊天室）──
            time.sleep(0.3)
            _si_click(_tw_x + int(_tw_w * 0.13), _tw_y + int(_tw_h * 0.12))
            time.sleep(0.2)
            results.append("🧹 已回到聊天列表")

            return "\n".join(results) if results else "Telegram導航完成"

        # ── 非 Telegram：通用視窗管理 ──
        # 先點一下目標螢幕中央，確保焦點在正確螢幕
        pyautogui.click(mon_cx, mon_cy)
        time.sleep(0.3)

        # 通用：把視窗拉到前景
        try:
            import win32gui, win32con
            def find_window(name):
                result = []
                win32gui.EnumWindows(
                    lambda h, _: result.append(h)
                    if name.lower() in win32gui.GetWindowText(h).lower() else None, None)
                return result[0] if result else None
            hw = find_window(app)
            if hw:
                win32gui.ShowWindow(hw, win32con.SW_RESTORE)
                win32gui.SetForegroundWindow(hw)
                time.sleep(0.6)
                results.append(f"✅ 切換到 {app} 視窗")
            else:
                results.append(f"⚠️ 找不到 {app} 視窗，嘗試開啟")
                pyautogui.hotkey("win", "s")
                time.sleep(0.5)
                pyautogui.write(app, interval=0.05)
                time.sleep(1)
                pyautogui.press("enter")
                time.sleep(2.5)
        except Exception as e:
            results.append(f"視窗切換：{e}")

        # 再次點目標螢幕確保焦點
        pyautogui.click(mon_cx, mon_cy)
        time.sleep(0.3)

        # LINE 專屬流程
        if "line" in app_lower:
            name_match = re.search(r'[「"](.*?)[」"]', task)
            name = name_match.group(1) if name_match else re.sub(r'(打開|找到|找|搜尋|對話|訊息)', '', task).strip()
            r = fetch_ocr_click(name, monitor)
            results.append(f"找對話 {name}：{r}")
            if input_text:
                time.sleep(0.5)
                pyautogui.write(input_text, interval=0.04)
                if "送出" in task or "回覆" in task:
                    pyautogui.press("enter")
                    results.append("📤 已送出")

        # 通用流程：用 OCR/視覺找目標
        else:
            r = fetch_ocr_click(task, monitor)
            results.append(f"OCR操作：{r}")
            if input_text:
                time.sleep(0.3)
                pyautogui.write(input_text, interval=0.04)
                results.append(f"已輸入：{input_text[:40]}")

        return "\n".join(results) if results else f"App導航完成：{app} / {task}"
    except Exception as e:
        return f"App導航失敗：{e}"


def fetch_wait_and_click(target_text: str, timeout: int = 15, monitor: int = 1, action_after: str = "click") -> str:
    try:
        import time
        start = time.time()
        interval = 1.0
        while time.time() - start < timeout:
            result = fetch_ocr_click(target_text, monitor, action_after if action_after != "none" else "click")
            if "✅" in result:
                return f"✅ 等待 {time.time()-start:.1f}s 後找到並點擊：{result}"
            if action_after == "none" and "找到" in result:
                return f"✅ 等待 {time.time()-start:.1f}s 後出現：{target_text}"
            time.sleep(interval)
        return f"⏰ 等待 {timeout}s 仍未出現「{target_text}」，超時"
    except Exception as e:
        return f"等待點擊失敗：{e}"


def fetch_drag_drop(from_x: int = None, from_y: int = None, to_x: int = None, to_y: int = None,
                    from_text: str = "", to_text: str = "", monitor: int = 1, duration: float = 0.5) -> str:
    try:
        import mss, ctypes as _c, ctypes.wintypes as _w, time as _t, re

        def get_abs(x, y):
            with mss.mss() as sct:
                mon = sct.monitors[monitor] if monitor < len(sct.monitors) else sct.monitors[1]
                return mon["left"] + x, mon["top"] + y

        # 起點
        if from_text:
            r = fetch_vision_locate(from_text, monitor, "locate_only")
            m = re.search(r'\((-?\d+), (-?\d+)\)', r)
            if m: fx, fy = int(m.group(1)), int(m.group(2))
            else: return f"拖曳起點找不到「{from_text}」：{r}"
        else:
            fx, fy = get_abs(from_x or 0, from_y or 0)

        # 終點
        if to_text:
            r2 = fetch_vision_locate(to_text, monitor, "locate_only")
            m2 = re.search(r'\((-?\d+), (-?\d+)\)', r2)
            if m2: tx, ty = int(m2.group(1)), int(m2.group(2))
            else: return f"拖曳終點找不到「{to_text}」：{r2}"
        else:
            tx, ty = get_abs(to_x or 0, to_y or 0)

        # SendInput 拖曳（支援螢幕2負座標）
        vl, vt, vw, vh, is_phys, sx, sy = _get_virtual_desktop()
        u32 = _c.windll.user32
        class MI(_c.Structure):
            _fields_ = [('dx',_w.LONG),('dy',_w.LONG),('mouseData',_w.DWORD),
                        ('dwFlags',_w.DWORD),('time',_w.DWORD),('dwExtraInfo',_c.POINTER(_c.c_ulong))]
        class U(_c.Union): _fields_ = [('mi', MI)]
        class INP(_c.Structure):
            _anonymous_ = ('u',); _fields_ = [('type',_w.DWORD),('u',U)]
        def _send(flags, dx=0, dy=0):
            i = INP(0, U(mi=MI(dx, dy, 0, flags, 0, None)))
            u32.SendInput(1, _c.byref(i), _c.sizeof(i))
        def _norm(lx, ly):
            px = round(lx * sx); py = round(ly * sy)
            return int((px-vl)*65535//vw), int((py-vt)*65535//vh)

        fnx, fny = _norm(fx, fy)
        tnx, tny = _norm(tx, ty)
        steps = max(10, int(duration * 60))
        _send(0x0001|0x8000|0x4000, fnx, fny); _t.sleep(0.1)
        _send(0x0002)  # LEFTDOWN
        for i in range(1, steps+1):
            ix = fnx + (tnx - fnx) * i // steps
            iy = fny + (tny - fny) * i // steps
            _send(0x0001|0x8000|0x4000, ix, iy); _t.sleep(duration / steps)
        _send(0x0004)  # LEFTUP
        return f"✅ 拖曳完成：({fx},{fy}) → ({tx},{ty})，耗時 {duration}s"
    except Exception as e:
        return f"拖曳失敗：{e}"


def fetch_read_screen(question: str = "描述螢幕上有什麼", monitor: int = 1) -> str:
    """截圖 → OCR + Vision → 回傳螢幕內容描述"""
    try:
        import anthropic, base64, io
        from PIL import Image as _PI
        img, _, _ = _cap_monitor_logical(monitor)
        if img.width > 2048:
            r = 2048 / img.width
            img = img.resize((2048, int(img.height*r)), _PI.LANCZOS)
        buf = io.BytesIO(); img.save(buf, format="JPEG", quality=92)
        if buf.tell() > 4 * 1024 * 1024:
            buf = io.BytesIO(); img.save(buf, format="JPEG", quality=80)
        b64 = base64.standard_b64encode(buf.getvalue()).decode()
        # OCR 輔助
        ocr_hint = ""
        try:
            import pytesseract
            ocr_text = pytesseract.image_to_string(img, lang="chi_tra+eng").strip()
            if ocr_text:
                ocr_hint = f"\n\nOCR偵測到的文字：{ocr_text[:500]}"
        except Exception:
            pass
        resp = anthropic.Anthropic().messages.create(
            model="claude-sonnet-4-6", max_tokens=1024,
            messages=[{"role": "user", "content": [
                {"type": "image", "source": {"type": "base64", "media_type": "image/jpeg", "data": b64}},
                {"type": "text", "text": f"這是螢幕{monitor}的截圖（{img.width}x{img.height}px）。{question}。請用繁體中文詳細描述畫面內容。{ocr_hint}"}
            ]}]
        )
        return f"📺 螢幕{monitor}內容：\n{resp.content[0].text}"
    except Exception as e:
        return f"讀取螢幕失敗：{e}"


def fetch_scroll_at(direction: str = "down", amount: int = 3,
                    x: int = None, y: int = None,
                    monitor: int = 1, description: str = "") -> str:
    """在指定位置滾動，支援所有螢幕包含螢幕2負座標"""
    try:
        import mss, time
        if description:
            img, ml, mt = _cap_monitor_logical(monitor)
            rx, ry = _vision_find(img, description)
            if rx is not None:
                abs_x, abs_y = ml + rx, mt + ry
            else:
                with mss.mss() as s:
                    m = s.monitors[monitor] if monitor < len(s.monitors) else s.monitors[1]
                    abs_x = m["left"] + m["width"] // 2
                    abs_y = m["top"] + m["height"] // 2
        elif x is not None and y is not None:
            with mss.mss() as s:
                m = s.monitors[monitor] if monitor < len(s.monitors) else s.monitors[1]
                abs_x = m["left"] + x; abs_y = m["top"] + y
        else:
            with mss.mss() as s:
                m = s.monitors[monitor] if monitor < len(s.monitors) else s.monitors[1]
                abs_x = m["left"] + m["width"] // 2
                abs_y = m["top"] + m["height"] // 2
        _si_scroll(abs_x, abs_y, amount, direction)
        return f"✅ 螢幕{monitor} 在({abs_x},{abs_y}) 向{direction}滾動 {amount} 格"
    except Exception as e:
        return f"滾動失敗：{e}"


def fetch_window_manager(action: str = "list", window_name: str = "") -> str:
    """視窗管理：列出所有視窗 / 切換焦點 / 最大化 / 最小化 / 關閉"""
    try:
        import win32gui, win32con, re
        results = []

        if action == "list":
            wins = []
            def _enum(h, _):
                if win32gui.IsWindowVisible(h):
                    t = win32gui.GetWindowText(h).strip()
                    if t: wins.append(f"[{h}] {t}")
                return True
            win32gui.EnumWindows(_enum, None)
            return "開啟中的視窗：\n" + "\n".join(wins[:30])

        # 找目標視窗
        matches = []
        def _find(h, _):
            if win32gui.IsWindowVisible(h):
                t = win32gui.GetWindowText(h)
                if window_name.lower() in t.lower():
                    matches.append(h)
            return True
        win32gui.EnumWindows(_find, None)

        if not matches:
            return f"找不到包含「{window_name}」的視窗"

        hw = matches[0]
        title = win32gui.GetWindowText(hw)

        if action == "focus":
            win32gui.ShowWindow(hw, win32con.SW_RESTORE)
            win32gui.SetForegroundWindow(hw)
            return f"✅ 已切換到視窗：{title}"
        elif action == "maximize":
            win32gui.ShowWindow(hw, win32con.SW_MAXIMIZE)
            return f"✅ 已最大化：{title}"
        elif action == "minimize":
            win32gui.ShowWindow(hw, win32con.SW_MINIMIZE)
            return f"✅ 已最小化：{title}"
        elif action == "close":
            win32gui.PostMessage(hw, win32con.WM_CLOSE, 0, 0)
            return f"✅ 已關閉：{title}"
        else:
            return f"未知動作：{action}，可用：list/focus/maximize/minimize/close"
    except Exception as e:
        return f"視窗管理失敗：{e}"


def fetch_crypto(coin: str, vs_currency: str = "usd") -> str:
    try:
        ticker_map = {
            "btc": "bitcoin", "eth": "ethereum", "sol": "solana",
            "bnb": "binancecoin", "xrp": "ripple", "doge": "dogecoin",
            "ada": "cardano", "dot": "polkadot", "matic": "matic-network",
            "avax": "avalanche-2", "link": "chainlink", "uni": "uniswap",
            "ltc": "litecoin", "bch": "bitcoin-cash", "atom": "cosmos",
            "trx": "tron", "etc": "ethereum-classic", "shib": "shiba-inu",
        }
        coin_id = ticker_map.get(coin.lower(), coin.lower())
        url = f"https://api.coingecko.com/api/v3/coins/{coin_id}?localization=false&tickers=false&community_data=false&developer_data=false"
        resp = requests.get(url, timeout=10)
        data = resp.json()
        if "error" in data:
            return f"找不到幣種「{coin}」，請確認名稱"
        md = data["market_data"]
        name = data["name"]
        sym = data["symbol"].upper()
        cur = vs_currency
        price = md["current_price"][cur]
        ch24 = md["price_change_percentage_24h"] or 0
        ch7d = md["price_change_percentage_7d"] or 0
        mcap = md["market_cap"][cur]
        vol = md["total_volume"][cur]
        hi24 = md["high_24h"][cur]
        lo24 = md["low_24h"][cur]
        ath = md["ath"][cur]
        ath_chg = md["ath_change_percentage"][cur] or 0
        cur_label = cur.upper()
        arrow = "▲" if ch24 >= 0 else "▼"
        mc_str = f"{mcap/1e12:.2f}T" if mcap >= 1e12 else f"{mcap/1e9:.2f}B"
        return (
            f"🪙 {name} ({sym})\n"
            f"💰 現價：{price:,.4f} {cur_label}  {arrow} {abs(ch24):.2f}% (24h)\n"
            f"📅 7日漲跌：{ch7d:+.2f}%\n"
            f"📊 24h 高低：{lo24:,.4f} ~ {hi24:,.4f}\n"
            f"💎 市值：{mc_str} {cur_label}\n"
            f"📦 24h 交易量：{vol/1e9:.2f}B {cur_label}\n"
            f"🏔 歷史高點：{ath:,.4f}（距高 {ath_chg:.1f}%）"
        )
    except Exception as e:
        return f"查詢加密幣「{coin}」失敗：{e}"


def fetch_forex(pair: str) -> str:
    try:
        import yfinance as yf
        ticker = yf.Ticker(pair)
        hist = ticker.history(period="5d")
        if hist.empty:
            return f"找不到匯率「{pair}」，請確認格式（如 USDTWD=X）"
        info = ticker.info
        current = hist["Close"].iloc[-1]
        prev = hist["Close"].iloc[-2] if len(hist) > 1 else current
        change = current - prev
        change_pct = (change / prev * 100) if prev else 0
        arrow = "▲" if change >= 0 else "▼"
        name = info.get("shortName", pair)
        return (
            f"💱 {name}\n"
            f"匯率：{current:.4f}  {arrow} {abs(change):.4f} ({change_pct:+.2f}%)\n"
            f"近5日高：{hist['High'].max():.4f}\n"
            f"近5日低：{hist['Low'].min():.4f}"
        )
    except Exception as e:
        return f"查詢匯率「{pair}」失敗：{e}"


def fetch_finance_news(source: str = "all", count: int = 5) -> str:
    try:
        import feedparser
        count = min(max(count, 1), 10)
        feeds = {
            "yahoo_tw": ("Yahoo奇摩財經", "https://tw.stock.yahoo.com/rss"),
            "cnyes":    ("鉅亨網",        "https://feeds.feedburner.com/cnyes"),
            "udn":      ("聯合財經網",    "https://udn.com/rssfeed/news/2/6644?ch=news"),
            "moneydj":  ("MoneyDJ",      "https://www.moneydj.com/KMDJ/RssNew/RssNewList.aspx?index=1&param="),
            "ctee":     ("工商時報",      "https://www.ctee.com.tw/feeds/latest"),
            "yahoo_us": ("Yahoo Finance", "https://finance.yahoo.com/news/rssindex"),
        }
        sources = list(feeds.keys()) if source == "all" else [source]
        results = []
        for src in sources:
            if src not in feeds:
                continue
            label, url = feeds[src]
            try:
                feed = feedparser.parse(url)
                items = feed.entries[:count]
                if not items:
                    results.append(f"📰 {label}：暫無資料")
                    continue
                lines = [f"📰 {label}"]
                for i, entry in enumerate(items, 1):
                    title = entry.get("title", "無標題")
                    lines.append(f"{i}. {title}")
                results.append("\n".join(lines))
            except Exception:
                results.append(f"📰 {label}：抓取失敗")
        return "\n\n".join(results) if results else "無法取得新聞"
    except Exception as e:
        return f"財經新聞失敗：{e}"


def ptt_search(keyword: str, board: str = "Gossiping", count: int = 5) -> str:
    try:
        from bs4 import BeautifulSoup
        import ssl, urllib3
        count = min(count, 10)
        # PTT 需要 session + 放寬 SSL + Cookie
        session = requests.Session()
        session.verify = False
        urllib3.disable_warnings(urllib3.exceptions.InsecureRequestWarning)
        headers = {
            "User-Agent": "Mozilla/5.0 (Windows NT 10.0; Win64; x64) AppleWebKit/537.36 Chrome/120.0",
            "Cookie": "over18=1",
            "Referer": "https://www.ptt.cc/",
        }
        search_url = f"https://www.ptt.cc/bbs/{board}/search?q={urllib.parse.quote(keyword)}"
        resp = session.get(search_url, headers=headers, timeout=12)
        soup = BeautifulSoup(resp.text, "html.parser")
        posts = soup.select(".r-ent")[:count]
        if not posts:
            # fallback: 用 Google News 搜尋 PTT 相關文章
            import feedparser
            fallback_url = f"https://news.google.com/rss/search?q=PTT+{urllib.parse.quote(keyword)}&hl=zh-Hant&gl=TW&ceid=TW:zh-Hant"
            feed = feedparser.parse(fallback_url)
            if feed.entries:
                lines = [f"📋 PTT 搜尋結果（Google News 補充）：{keyword}\n"]
                for i, e in enumerate(feed.entries[:count], 1):
                    lines.append(f"{i}. {e.get('title', '')}")
                return "\n".join(lines)
            return f"PTT {board} 版找不到「{keyword}」相關文章"
        lines = [f"📋 PTT/{board} 搜尋：{keyword}\n"]
        for post in posts:
            title_el = post.select_one(".title a")
            meta_el = post.select_one(".meta .author")
            date_el = post.select_one(".meta .date")
            nrec_el = post.select_one(".nrec span")
            if not title_el:
                continue
            title = title_el.get_text(strip=True)
            author = meta_el.get_text(strip=True) if meta_el else ""
            date = date_el.get_text(strip=True) if date_el else ""
            nrec = nrec_el.get_text(strip=True) if nrec_el else "0"
            post_url = "https://www.ptt.cc" + title_el["href"]
            lines.append(f"🗂 {title}")
            lines.append(f"   推文：{nrec}　作者：{author}　{date}")
            try:
                p_resp = session.get(post_url, headers=headers, timeout=6)
                p_soup = BeautifulSoup(p_resp.text, "html.parser")
                content_el = p_soup.select_one("#main-content")
                if content_el:
                    for tag in content_el.select(".f2, .push, #article-polling"):
                        tag.decompose()
                    content = content_el.get_text(separator=" ", strip=True)[:200]
                    lines.append(f"   摘要：{content}")
            except Exception:
                pass
        return "\n".join(lines)
    except Exception as e:
        return f"PTT 搜尋失敗：{e}"


def multi_perspective(topic: str, lang: str = "zh-tw") -> str:
    try:
        from ddgs import DDGS
        results = {}
        queries = {
            "支持/正面觀點": f"{topic} 優點 支持 正面",
            "反對/批評觀點": f"{topic} 缺點 反對 批評 問題",
            "中立/分析觀點": f"{topic} 分析 評估 影響 研究",
        }
        if lang == "en":
            queries = {
                "Pro / Positive": f"{topic} benefits support positive",
                "Con / Critical": f"{topic} criticism problems negative against",
                "Neutral / Analysis": f"{topic} analysis impact research objective",
            }
        with DDGS() as ddgs:
            for label, q in queries.items():
                items = list(ddgs.text(q, region=lang, max_results=3))
                lines = [f"── {label} ──"]
                for r in items:
                    lines.append(f"• {r['title']}")
                    lines.append(f"  {r['body'][:150]}")
                results[label] = "\n".join(lines)
        output = [f"🔍 多角度分析：{topic}\n"]
        output.extend(results.values())
        return "\n\n".join(output)
    except Exception as e:
        return f"多角度分析失敗：{e}"


def fetch_google_trends(keywords: list, timeframe: str = "today 3-m", geo: str = "TW") -> str:
    try:
        from pytrends.request import TrendReq
        keywords = keywords[:5]
        pt = TrendReq(hl="zh-TW", tz=480, timeout=(5, 15))
        pt.build_payload(keywords, timeframe=timeframe, geo=geo)
        df = pt.interest_over_time()
        if df.empty:
            return f"找不到「{', '.join(keywords)}」的趨勢資料"
        lines = [f"📈 Google Trends（{geo}，{timeframe}）\n"]
        for kw in keywords:
            if kw not in df.columns:
                continue
            avg = df[kw].mean()
            peak = df[kw].max()
            peak_date = str(df[kw].idxmax())[:10]
            recent = df[kw].iloc[-4:].mean()
            trend = "上升 📈" if df[kw].iloc[-1] > df[kw].iloc[-8] else "下降 📉"
            lines.append(
                f"🔍 {kw}\n"
                f"   平均熱度：{avg:.0f}　峰值：{peak}（{peak_date}）\n"
                f"   近期熱度：{recent:.0f}　趨勢：{trend}"
            )
        # 相關搜尋
        try:
            related = pt.related_queries()
            for kw in keywords[:2]:
                if kw in related and related[kw]["top"] is not None:
                    top_q = related[kw]["top"]["query"].head(3).tolist()
                    lines.append(f"\n「{kw}」相關搜尋：{', '.join(top_q)}")
        except Exception:
            pass
        return "\n".join(lines)
    except Exception as e:
        return f"Google Trends 查詢失敗：{e}"


def read_webpage(url: str, max_chars: int = 3000) -> str:
    """讀取網頁內容，先用 requests，內容太少或失敗時自動 fallback 到 Playwright"""
    def _parse_html(html, url):
        from bs4 import BeautifulSoup
        soup = BeautifulSoup(html, "html.parser")
        for tag in soup(["script", "style", "nav", "footer", "header", "aside", "ads"]):
            tag.decompose()
        content_el = soup.find("article") or soup.find("main") or soup.find("body")
        text = content_el.get_text(separator="\n", strip=True) if content_el else soup.get_text(separator="\n", strip=True)
        lines = [l.strip() for l in text.splitlines() if l.strip()]
        clean = "\n".join(lines)
        title = soup.title.string.strip() if soup.title else url
        return title, clean

    # 策略1：requests（快）
    try:
        headers = {"User-Agent": "Mozilla/5.0 (Windows NT 10.0; Win64; x64) AppleWebKit/537.36"}
        resp = requests.get(url, headers=headers, timeout=12)
        resp.encoding = resp.apparent_encoding
        title, clean = _parse_html(resp.text, url)
        # 如果內容太少（<100字），可能是動態渲染頁面，fallback 到 Playwright
        if len(clean) >= 100:
            result = f"【{title}】\n{url}\n\n{clean[:max_chars]}"
            if len(clean) > max_chars:
                result += f"\n\n（內容已截斷，共 {len(clean)} 字）"
            return result
    except Exception:
        pass

    # 策略2：Playwright headless（能渲染 JavaScript 動態頁面）
    try:
        import subprocess, sys, json as _json_rw, textwrap
        _script = textwrap.dedent(f"""
import sys, json
from playwright.sync_api import sync_playwright
with sync_playwright() as p:
    b = p.chromium.launch(headless=True)
    page = b.new_page(user_agent="Mozilla/5.0 (Windows NT 10.0; Win64; x64) AppleWebKit/537.36")
    page.goto({_json_rw.dumps(url)}, wait_until="networkidle", timeout=20000)
    page.wait_for_timeout(2000)
    html = page.content()
    b.close()
    print(html)
""")
        proc = subprocess.run(
            [sys.executable, "-c", _script],
            capture_output=True, text=True, encoding="utf-8", errors="replace",
            timeout=30, env={**__import__("os").environ, "PYTHONIOENCODING": "utf-8"}
        )
        if proc.returncode == 0 and proc.stdout.strip():
            title, clean = _parse_html(proc.stdout, url)
            if clean:
                result = f"【{title}】（Playwright）\n{url}\n\n{clean[:max_chars]}"
                if len(clean) > max_chars:
                    result += f"\n\n（內容已截斷，共 {len(clean)} 字）"
                return result
    except Exception:
        pass

    return f"網頁讀取失敗：無法透過 requests 或 Playwright 取得 {url} 的內容"


def wikipedia_search(query: str, lang: str = "zh") -> str:
    try:
        search_url = f"https://{lang}.wikipedia.org/w/api.php"
        # 先搜尋
        params = {"action": "search", "list": "search", "srsearch": query,
                  "format": "json", "srlimit": 1}
        resp = requests.get(search_url, params=params, timeout=10)
        results = resp.json().get("query", {}).get("search", [])
        if not results:
            return f"Wikipedia 找不到「{query}」相關條目"
        title = results[0]["title"]
        # 取得摘要
        params2 = {"action": "query", "titles": title, "prop": "extracts",
                   "exintro": True, "explaintext": True, "format": "json"}
        resp2 = requests.get(search_url, params=params2, timeout=10)
        pages = resp2.json().get("query", {}).get("pages", {})
        page = next(iter(pages.values()))
        extract = page.get("extract", "無法取得內容")
        extract = extract[:2500]
        return f"📖 Wikipedia：{title}\n\n{extract}"
    except Exception as e:
        return f"Wikipedia 查詢失敗：{e}"


def search_news(query: str, lang: str = "zh-TW", count: int = 6) -> str:
    try:
        import feedparser
        count = min(count, 10)
        lang_map = {"zh-TW": "zh-Hant&gl=TW&ceid=TW:zh-Hant",
                    "zh-CN": "zh-Hans&gl=CN&ceid=CN:zh-Hans",
                    "en-US": "en&gl=US&ceid=US:en"}
        param = lang_map.get(lang, lang_map["zh-TW"])
        url = f"https://news.google.com/rss/search?q={urllib.parse.quote(query)}&hl={param}"
        feed = feedparser.parse(url)
        if not feed.entries:
            return f"找不到「{query}」的新聞"
        lines = [f"📰 {query} 最新新聞\n"]
        for i, entry in enumerate(feed.entries[:count], 1):
            title = entry.get("title", "無標題")
            pub = entry.get("published", "")[:16]
            lines.append(f"{i}. {title}（{pub}）")
        return "\n".join(lines)
    except Exception as e:
        return f"新聞搜尋失敗：{e}"


def youtube_summary(url: str) -> str:
    try:
        from youtube_transcript_api import YouTubeTranscriptApi
        import re
        # 從 URL 提取 video ID
        patterns = [
            r"(?:v=|youtu\.be/|/embed/|/v/)([a-zA-Z0-9_-]{11})",
        ]
        video_id = url if re.match(r"^[a-zA-Z0-9_-]{11}$", url) else None
        if not video_id:
            for pat in patterns:
                m = re.search(pat, url)
                if m:
                    video_id = m.group(1)
                    break
        if not video_id:
            return f"無法從網址提取 YouTube 影片 ID：{url}"
        # 嘗試取得字幕（優先繁中→簡中→英文）
        transcript = None
        for lang in [["zh-TW", "zh-Hant"], ["zh", "zh-Hans"], ["en"]]:
            try:
                transcript = YouTubeTranscriptApi.get_transcript(video_id, languages=lang)
                break
            except Exception:
                continue
        if transcript is None:
            try:
                transcript = YouTubeTranscriptApi.get_transcript(video_id)
            except Exception as e:
                return f"無法取得字幕（影片可能無字幕或受限）：{e}"
        # 合併字幕文字
        full_text = " ".join(t["text"] for t in transcript)
        # 截取前 3000 字
        summary = full_text[:3000]
        if len(full_text) > 3000:
            summary += f"\n\n（字幕共 {len(full_text)} 字，已截取前段）"
        return f"🎬 YouTube 字幕摘要\n影片ID：{video_id}\n\n{summary}"
    except Exception as e:
        return f"YouTube 字幕擷取失敗：{e}"


def analyze_pdf(path: str, max_chars: int = 4000) -> str:
    try:
        import pdfplumber
        with pdfplumber.open(path) as pdf:
            total_pages = len(pdf.pages)
            texts = []
            for page in pdf.pages[:20]:  # 最多讀前 20 頁
                t = page.extract_text()
                if t:
                    texts.append(t)
        full_text = "\n".join(texts)
        if not full_text.strip():
            return "PDF 無法提取文字（可能是掃描圖片 PDF）"
        result = f"📄 PDF 分析（共 {total_pages} 頁）\n\n{full_text[:max_chars]}"
        if len(full_text) > max_chars:
            result += f"\n\n（內容已截斷，共 {len(full_text)} 字）"
        return result
    except Exception as e:
        return f"PDF 讀取失敗：{e}"


def execute_ddg_search(query: str, region: str = "zh-tw", max_results: int = 5) -> str:
    try:
        from ddgs import DDGS
        max_results = min(max(max_results, 1), 10)
        results = []
        with DDGS() as ddgs:
            for r in ddgs.text(query, region=region, max_results=max_results):
                title = r.get("title", "")
                body = r.get("body", "")
                href = r.get("href", "")
                # 跳過大陸域名（.cn）避免超時
                if ".cn/" in href and "taiwan" not in href.lower():
                    pass
                results.append(f"🔍 {title}\n   {body[:120]}\n   {href}")
        return "\n\n".join(results) if results else "無搜尋結果"
    except Exception as e:
        return f"DuckDuckGo 搜尋失敗：{e}"


def fetch_stock_advanced(symbol: str, indicators: list = None) -> str:
    try:
        import yfinance as yf
        import ta as ta_lib
        if indicators is None:
            indicators = ["macd", "bb", "kd"]
        ticker = yf.Ticker(symbol)
        hist = ticker.history(period="6mo")
        if hist.empty:
            return f"找不到「{symbol}」的資料"
        name = ticker.info.get("longName") or ticker.info.get("shortName") or symbol
        close = hist["Close"]
        high = hist["High"]
        low = hist["Low"]
        current = close.iloc[-1]
        result = [f"📊 {name} ({symbol}) 進階技術分析\n現價：{current:.2f}\n"]
        if "macd" in indicators:
            macd_ind = ta_lib.trend.MACD(close)
            macd_v = macd_ind.macd().iloc[-1]
            signal_v = macd_ind.macd_signal().iloc[-1]
            hist_v = macd_ind.macd_diff().iloc[-1]
            cross = "金叉 📈" if macd_v > signal_v else "死叉 📉"
            result.append(
                f"── MACD ──\n"
                f"MACD：{macd_v:.3f}　Signal：{signal_v:.3f}　Histogram：{hist_v:.3f}\n"
                f"狀態：{cross}"
            )
        if "bb" in indicators:
            bb = ta_lib.volatility.BollingerBands(close)
            upper = bb.bollinger_hband().iloc[-1]
            mid = bb.bollinger_mavg().iloc[-1]
            lower = bb.bollinger_lband().iloc[-1]
            width = (upper - lower) / mid * 100 if mid != 0 else 0
            if current >= upper:
                bb_pos = "觸上軌（超買警示）⚠️"
            elif current <= lower:
                bb_pos = "觸下軌（超賣機會）💡"
            else:
                pos_pct = (current - lower) / (upper - lower) * 100
                bb_pos = f"通道內 {pos_pct:.0f}%"
            result.append(
                f"\n── 布林通道（BB）──\n"
                f"上軌：{upper:.2f}　中軌：{mid:.2f}　下軌：{lower:.2f}\n"
                f"帶寬：{width:.1f}%　位置：{bb_pos}"
            )
        if "kd" in indicators:
            stoch = ta_lib.momentum.StochasticOscillator(high, low, close)
            k = stoch.stoch().iloc[-1]
            d = stoch.stoch_signal().iloc[-1]
            if k > 80:
                kd_note = "超買區（K>80）⚠️"
            elif k < 20:
                kd_note = "超賣區（K<20）💡"
            else:
                kd_note = "中性區間"
            cross_kd = "K上穿D（買訊）📈" if k > d else "K下穿D（賣訊）📉"
            result.append(
                f"\n── KD 指標 ──\n"
                f"K：{k:.1f}　D：{d:.1f}\n"
                f"狀態：{kd_note}　交叉：{cross_kd}"
            )
        return "\n".join(result)
    except Exception as e:
        return f"進階技術分析失敗：{e}"


def fetch_fundamentals(symbol: str) -> str:
    try:
        import yfinance as yf
        ticker = yf.Ticker(symbol)
        info = ticker.info
        name = info.get("longName") or info.get("shortName") or symbol
        currency = info.get("currency", "")

        # 基本面指標
        roe = info.get("returnOnEquity")
        roa = info.get("returnOnAssets")
        profit_margin = info.get("profitMargins")
        gross_margin = info.get("grossMargins")
        operating_margin = info.get("operatingMargins")
        debt_equity = info.get("debtToEquity")
        current_ratio = info.get("currentRatio")
        eps = info.get("trailingEps")
        eps_fwd = info.get("forwardEps")
        pe = info.get("trailingPE")
        pe_fwd = info.get("forwardPE")
        pb = info.get("priceToBook")
        rev_growth = info.get("revenueGrowth")
        earn_growth = info.get("earningsGrowth")
        dividend_yield = info.get("dividendYield")

        # 分析師評級
        target_mean = info.get("targetMeanPrice")
        target_high = info.get("targetHighPrice")
        target_low = info.get("targetLowPrice")
        recommend = info.get("recommendationKey", "")
        recommend_map = {"strong_buy": "強力買進 💚", "buy": "買進 🟢", "hold": "持有 🟡",
                         "sell": "賣出 🔴", "strong_sell": "強力賣出 ❌"}
        recommend_str = recommend_map.get(recommend, recommend)
        num_analysts = info.get("numberOfAnalystOpinions", 0)

        lines = [f"📋 {name} ({symbol}) 深度基本面\n"]

        lines.append("── 獲利能力 ──")
        if roe: lines.append(f"ROE（股東權益報酬）：{roe*100:.1f}%")
        if roa: lines.append(f"ROA（資產報酬）：{roa*100:.1f}%")
        if gross_margin: lines.append(f"毛利率：{gross_margin*100:.1f}%")
        if operating_margin: lines.append(f"營業利益率：{operating_margin*100:.1f}%")
        if profit_margin: lines.append(f"淨利率：{profit_margin*100:.1f}%")

        lines.append("\n── 估值 ──")
        if pe: lines.append(f"本益比（P/E）：{pe:.1f}")
        if pe_fwd: lines.append(f"預估本益比：{pe_fwd:.1f}")
        if pb: lines.append(f"股價淨值比（P/B）：{pb:.2f}")
        if eps: lines.append(f"EPS（過去12月）：{eps:.2f} {currency}")
        if eps_fwd: lines.append(f"EPS 預估：{eps_fwd:.2f} {currency}")

        lines.append("\n── 成長性 ──")
        if rev_growth: lines.append(f"營收年增率：{rev_growth*100:+.1f}%")
        if earn_growth: lines.append(f"獲利年增率：{earn_growth*100:+.1f}%")

        lines.append("\n── 財務健康 ──")
        if debt_equity: lines.append(f"負債股權比：{debt_equity:.1f}%")
        if current_ratio: lines.append(f"流動比率：{current_ratio:.2f}")
        if dividend_yield: lines.append(f"殖利率：{dividend_yield*100:.2f}%")

        if target_mean and num_analysts:
            lines.append(f"\n── 分析師（{num_analysts} 位）──")
            lines.append(f"評級：{recommend_str}")
            lines.append(f"目標價：{target_low:.2f} ~ {target_high:.2f}（均值 {target_mean:.2f} {currency}）")

        return "\n".join(lines)
    except Exception as e:
        return f"基本面查詢失敗：{e}"


def fetch_market_sentiment() -> str:
    try:
        import yfinance as yf
        # VIX
        vix_hist = yf.Ticker("^VIX").history(period="5d")
        vix = vix_hist["Close"].iloc[-1] if not vix_hist.empty else None
        vix_prev = vix_hist["Close"].iloc[-2] if len(vix_hist) > 1 else vix

        if vix:
            if vix >= 40: vix_note = "極度恐慌 😱"
            elif vix >= 30: vix_note = "高度恐慌 😰"
            elif vix >= 20: vix_note = "輕度緊張 😟"
            elif vix >= 12: vix_note = "正常 😐"
            else: vix_note = "市場過度樂觀 😎"

        # Fear & Greed Index (CNN)
        try:
            fg_resp = requests.get(
                "https://production.dataviz.cnn.io/index/fearandgreed/graphdata",
                headers={"User-Agent": "Mozilla/5.0"}, timeout=8
            )
            fg_data = fg_resp.json()
            fg_score = fg_data["fear_and_greed"]["score"]
            fg_rating = fg_data["fear_and_greed"]["rating"]
            fg_prev = fg_data["fear_and_greed"]["previous_close"]
            fg_str = f"{fg_score:.0f}/100 — {fg_rating}"
            fg_change = fg_score - fg_prev
        except Exception:
            fg_str = "無法取得"
            fg_change = 0

        # S&P500 and Nasdaq 當日走勢
        sp_hist = yf.Ticker("^GSPC").history(period="5d")
        ndx_hist = yf.Ticker("^IXIC").history(period="5d")
        sp_chg = ((sp_hist["Close"].iloc[-1] / sp_hist["Close"].iloc[-2]) - 1) * 100 if len(sp_hist) > 1 else 0
        ndx_chg = ((ndx_hist["Close"].iloc[-1] / ndx_hist["Close"].iloc[-2]) - 1) * 100 if len(ndx_hist) > 1 else 0

        lines = ["🌡 市場情緒儀表板\n"]
        if vix:
            vix_chg = vix - vix_prev if vix_prev else 0
            lines.append(f"VIX 波動率：{vix:.2f}（{vix_note}）{vix_chg:+.2f}")
        lines.append(f"恐慌貪婪指數：{fg_str}（{fg_change:+.1f}）")
        lines.append(f"\nS&P 500：{sp_chg:+.2f}%")
        lines.append(f"Nasdaq：{ndx_chg:+.2f}%")

        return "\n".join(lines)
    except Exception as e:
        return f"市場情緒查詢失敗：{e}"


def fetch_etf(symbol: str) -> str:
    try:
        import yfinance as yf
        ticker = yf.Ticker(symbol)
        info = ticker.info
        name = info.get("longName") or info.get("shortName") or symbol
        currency = info.get("currency", "")

        hist = ticker.history(period="1y")
        if hist.empty:
            return f"找不到 ETF「{symbol}」"

        current = hist["Close"].iloc[-1]
        ret_1m = ((hist["Close"].iloc[-1] / hist["Close"].iloc[-22]) - 1) * 100 if len(hist) > 22 else None
        ret_3m = ((hist["Close"].iloc[-1] / hist["Close"].iloc[-66]) - 1) * 100 if len(hist) > 66 else None
        ret_1y = ((hist["Close"].iloc[-1] / hist["Close"].iloc[0]) - 1) * 100 if len(hist) > 5 else None

        expense = info.get("annualReportExpenseRatio")
        div_yield = info.get("dividendYield") or info.get("yield")
        total_assets = info.get("totalAssets")
        category = info.get("category", "")

        # 前10大持股
        try:
            holdings = ticker.funds_data.top_holdings
            top_str = ""
            if holdings is not None and not holdings.empty:
                top_items = []
                for _, row in holdings.head(5).iterrows():
                    n = row.get("Name") or row.get("name") or ""
                    w = row.get("Holding Percent") or row.get("holdingPercent") or row.get("weight") or 0
                    top_items.append(f"  {n}（{w*100:.1f}%）")
                if top_items:
                    top_str = "\n前5大持股：\n" + "\n".join(top_items)
        except Exception:
            top_str = ""

        lines = [f"📦 {name} ({symbol})\n現價：{current:.2f} {currency}\n"]
        if category: lines.append(f"類型：{category}")
        if total_assets: lines.append(f"規模：{total_assets/1e9:.1f}B {currency}")
        if expense: lines.append(f"費用率：{expense*100:.2f}%")
        if div_yield: lines.append(f"配息殖利率：{div_yield*100:.2f}%")
        lines.append("\n── 績效 ──")
        if ret_1m: lines.append(f"近1月：{ret_1m:+.2f}%")
        if ret_3m: lines.append(f"近3月：{ret_3m:+.2f}%")
        if ret_1y: lines.append(f"近1年：{ret_1y:+.2f}%")
        if top_str: lines.append(top_str)

        return "\n".join(lines)
    except Exception as e:
        return f"ETF 查詢失敗：{e}"


def fetch_earnings(symbol: str) -> str:
    try:
        import yfinance as yf
        ticker = yf.Ticker(symbol)
        info = ticker.info
        name = info.get("longName") or info.get("shortName") or symbol

        lines = [f"📊 {name} ({symbol}) 財報趨勢\n"]

        # 季度財報（用 quarterly_income_stmt，quarterly_earnings 已棄用）
        try:
            qi = ticker.quarterly_income_stmt
            if qi is not None and not qi.empty and "Basic EPS" in qi.index:
                lines.append("── 近幾季 EPS ──")
                eps_row = qi.loc["Basic EPS"]
                for col in eps_row.index[:6]:
                    val = eps_row[col]
                    if val is not None and str(val) != "nan":
                        quarter = str(col)[:10]
                        lines.append(f"  {quarter}：EPS {float(val):.2f}")
        except Exception:
            pass

        # 年度財報趨勢
        try:
            fin = ticker.financials
            if fin is not None and not fin.empty:
                lines.append("\n── 年度財務 ──")
                rev_row = fin.loc["Total Revenue"] if "Total Revenue" in fin.index else None
                ni_row = fin.loc["Net Income"] if "Net Income" in fin.index else None
                cols = fin.columns[:4]
                if rev_row is not None:
                    vals = [f"{rev_row[c]/1e9:.1f}B" for c in cols if c in rev_row.index]
                    lines.append(f"  營收：{' → '.join(vals)}")
                if ni_row is not None:
                    vals = [f"{ni_row[c]/1e9:.1f}B" for c in cols if c in ni_row.index]
                    lines.append(f"  淨利：{' → '.join(vals)}")
        except Exception:
            pass

        # 下次財報日
        try:
            cal = ticker.calendar
            if cal is not None:
                ed = cal.get("Earnings Date") or cal.get("earningsDate")
                if ed is not None:
                    if hasattr(ed, '__iter__'):
                        ed = list(ed)[0]
                    lines.append(f"\n下次財報日：{str(ed)[:10]}")
        except Exception:
            pass

        return "\n".join(lines) if len(lines) > 1 else f"找不到 {symbol} 的財報資料"
    except Exception as e:
        return f"財報查詢失敗：{e}"


def generate_candlestick(symbol: str, period: str = "3mo") -> tuple[bytes | None, str]:
    """回傳 (PNG bytes, 型態分析文字)"""
    try:
        import yfinance as yf
        import mplfinance as mpf
        import tempfile, ta as ta_lib
        ticker = yf.Ticker(symbol)
        hist = ticker.history(period=period)
        if hist.empty:
            return None, f"找不到「{symbol}」資料"

        name = ticker.info.get("shortName") or symbol
        close = hist["Close"]
        ma20 = close.rolling(20).mean()
        ma60 = close.rolling(60).mean()

        add_plots = [
            mpf.make_addplot(ma20, color="orange", width=1.2, label="MA20"),
            mpf.make_addplot(ma60, color="purple", width=1.2, label="MA60"),
        ]

        with tempfile.NamedTemporaryFile(suffix=".png", delete=False) as f:
            out_path = f.name

        mpf.plot(
            hist, type="candle", style="charles",
            title=f"{name} ({symbol})",
            ylabel="Price", volume=True,
            addplot=add_plots,
            savefig=dict(fname=out_path, dpi=150, bbox_inches="tight"),
            figsize=(12, 7)
        )

        # 型態辨識
        patterns = []
        c = close.values
        n = len(c)
        if n >= 20:
            recent = c[-20:]
            peak_idx = recent.argmax()
            trough_idx = recent.argmin()
            # 簡單型態判斷
            if c[-1] > ma20.iloc[-1] > ma60.iloc[-1]:
                patterns.append("多頭排列（MA20>MA60，趨勢向上）📈")
            elif c[-1] < ma20.iloc[-1] < ma60.iloc[-1]:
                patterns.append("空頭排列（MA20<MA60，趨勢向下）📉")
            if peak_idx < 5 and c[-1] < recent[peak_idx] * 0.95:
                patterns.append("近期高點已過，回落中")
            if trough_idx < 5 and c[-1] > recent[trough_idx] * 1.05:
                patterns.append("近期低點反彈，留意支撐")
            # 突破判斷
            resistance = max(c[-20:-5])
            support = min(c[-20:-5])
            if c[-1] > resistance:
                patterns.append(f"突破近期壓力 {resistance:.2f} ⚡")
            elif c[-1] < support:
                patterns.append(f"跌破近期支撐 {support:.2f} ⚠️")

        pattern_str = "\n".join(patterns) if patterns else "無明顯型態訊號"

        with open(out_path, "rb") as f:
            img_bytes = f.read()
        Path(out_path).unlink(missing_ok=True)
        return img_bytes, pattern_str
    except Exception as e:
        return None, f"K線圖生成失敗：{e}"


def compare_stocks(symbols: list, metrics: list = None) -> str:
    try:
        import yfinance as yf
        if metrics is None or "all" in metrics:
            metrics = ["price", "pe", "roe", "margin", "growth"]

        rows = []
        for sym in symbols[:5]:
            try:
                info = yf.Ticker(sym).info
                hist = yf.Ticker(sym).history(period="1mo")
                ret_1m = ((hist["Close"].iloc[-1] / hist["Close"].iloc[0]) - 1) * 100 if len(hist) > 1 else None
                row = {
                    "symbol": sym,
                    "name": (info.get("shortName") or sym)[:15],
                    "price": info.get("currentPrice") or info.get("regularMarketPrice"),
                    "ret_1m": ret_1m,
                    "pe": info.get("trailingPE"),
                    "pb": info.get("priceToBook"),
                    "roe": info.get("returnOnEquity"),
                    "margin": info.get("profitMargins"),
                    "rev_growth": info.get("revenueGrowth"),
                    "earn_growth": info.get("earningsGrowth"),
                    "div_yield": info.get("dividendYield"),
                    "mkt_cap": info.get("marketCap"),
                }
                rows.append(row)
            except Exception:
                rows.append({"symbol": sym, "name": sym})

        if not rows:
            return "無法取得比較資料"

        lines = [f"📊 股票比較：{' vs '.join(symbols[:5])}\n"]

        def fmt(v, pct=False, mult=100):
            if v is None: return "N/A"
            if pct: return f"{v*mult:+.1f}%"
            return f"{v:.2f}"

        for r in rows:
            sym = r["symbol"]
            name = r.get("name", sym)
            mc = r.get("mkt_cap")
            mc_str = f"{mc/1e12:.2f}T" if mc and mc >= 1e12 else (f"{mc/1e9:.1f}B" if mc else "N/A")
            lines.append(f"── {name} ({sym}) ──")
            if "price" in metrics and r.get("price"):
                lines.append(f"  現價：{r['price']:.2f}　近1月：{fmt(r.get('ret_1m'), False)+'%' if r.get('ret_1m') else 'N/A'}")
            lines.append(f"  市值：{mc_str}")
            if "pe" in metrics:
                lines.append(f"  P/E：{fmt(r.get('pe'))}　P/B：{fmt(r.get('pb'))}")
            if "roe" in metrics and r.get("roe"):
                lines.append(f"  ROE：{r['roe']*100:.1f}%")
            if "margin" in metrics and r.get("margin"):
                lines.append(f"  淨利率：{r['margin']*100:.1f}%")
            if "growth" in metrics:
                lines.append(f"  營收成長：{fmt(r.get('rev_growth'), True)}　獲利成長：{fmt(r.get('earn_growth'), True)}")
            if r.get("div_yield"):
                lines.append(f"  殖利率：{r['div_yield']*100:.2f}%")

        return "\n".join(lines)
    except Exception as e:
        return f"股票比較失敗：{e}"


def fetch_macro(indicator: str) -> str:
    try:
        if indicator == "fed_rate":
            import yfinance as yf
            ticker = yf.Ticker("^IRX")
            hist = ticker.history(period="1mo")
            if not hist.empty:
                rate = hist["Close"].iloc[-1]
                prev = hist["Close"].iloc[-2] if len(hist) > 1 else rate
                change = rate - prev
                return (
                    f"🏦 美國短期利率（13週國庫券）\n"
                    f"當前：{rate:.2f}%　前日：{prev:.2f}%　變化：{change:+.2f}%\n"
                    f"近1月高低：{hist['Low'].min():.2f}% ~ {hist['High'].max():.2f}%\n"
                    f"（聯邦基金目標利率請參考 federalreserve.gov）"
                )
            return "無法取得利率資料"
        if indicator == "nonfarm":
            return (
                "📊 美國非農就業（Non-Farm Payrolls）\n"
                "每月第一個週五由美國勞工部公布。\n"
                "建議用 get_finance_news 工具搜尋最新數據，或查詢 bls.gov。"
            )
        wb_map = {
            "cpi":          ("美國 CPI 通膨年增率",  "FP.CPI.TOTL.ZG",  "US"),
            "unemployment": ("美國失業率",            "SL.UEM.TOTL.ZS",  "US"),
            "gdp":          ("美國 GDP 年增率",       "NY.GDP.MKTP.KD.ZG","US"),
        }
        label, wb_code, country = wb_map[indicator]
        url = f"https://api.worldbank.org/v2/country/{country}/indicator/{wb_code}?format=json&mrv=5"
        resp = requests.get(url, timeout=10)
        data = resp.json()
        if len(data) < 2 or not data[1]:
            return f"無法取得 {label} 資料"
        entries = [e for e in data[1] if e.get("value") is not None][:5]
        lines = [f"📈 {label}（世界銀行）"]
        for e in entries:
            lines.append(f"{e['date']}：{e['value']:.2f}%")
        return "\n".join(lines)
    except Exception as e:
        return f"總經指標查詢失敗：{e}"


_PORTFOLIO_DB = Path(__file__).parent / "portfolio.db"

def _init_portfolio_db():
    conn = sqlite3.connect(_PORTFOLIO_DB)
    conn.execute("""
        CREATE TABLE IF NOT EXISTS portfolio (
            id INTEGER PRIMARY KEY AUTOINCREMENT,
            chat_id INTEGER NOT NULL,
            symbol TEXT NOT NULL,
            shares REAL NOT NULL,
            cost REAL NOT NULL,
            created_at TIMESTAMP DEFAULT CURRENT_TIMESTAMP
        )
    """)
    conn.commit()
    conn.close()

def execute_portfolio(action: str, chat_id: int = 0, symbol: str = "",
                      shares: float = 0, cost: float = 0) -> str:
    try:
        _init_portfolio_db()
        if action == "add":
            conn = sqlite3.connect(_PORTFOLIO_DB)
            existing = conn.execute(
                "SELECT id, shares, cost FROM portfolio WHERE chat_id=? AND symbol=?",
                (chat_id, symbol.upper())
            ).fetchone()
            if existing:
                new_shares = existing[1] + shares
                new_cost = (existing[1] * existing[2] + shares * cost) / new_shares
                conn.execute("UPDATE portfolio SET shares=?, cost=? WHERE id=?",
                             (new_shares, new_cost, existing[0]))
            else:
                conn.execute(
                    "INSERT INTO portfolio (chat_id, symbol, shares, cost) VALUES (?, ?, ?, ?)",
                    (chat_id, symbol.upper(), shares, cost)
                )
            conn.commit(); conn.close()
            return f"✅ 已新增 {symbol.upper()} {shares} 股，成本 {cost:.2f}"
        elif action == "remove":
            conn = sqlite3.connect(_PORTFOLIO_DB)
            conn.execute("DELETE FROM portfolio WHERE chat_id=? AND symbol=?",
                         (chat_id, symbol.upper()))
            conn.commit(); conn.close()
            return f"✅ 已移除 {symbol.upper()}"
        elif action == "clear":
            conn = sqlite3.connect(_PORTFOLIO_DB)
            conn.execute("DELETE FROM portfolio WHERE chat_id=?", (chat_id,))
            conn.commit(); conn.close()
            return "✅ 已清空投資組合"
        elif action == "view":
            conn = sqlite3.connect(_PORTFOLIO_DB)
            rows = conn.execute(
                "SELECT symbol, shares, cost FROM portfolio WHERE chat_id=?", (chat_id,)
            ).fetchall()
            conn.close()
            if not rows:
                return "目前沒有持股紀錄，用 add 新增持股"
            import yfinance as yf
            lines = ["📁 投資組合\n"]
            total_cost = total_value = 0
            for sym, sh, ct in rows:
                try:
                    h = yf.Ticker(sym).history(period="1d")
                    cur_price = h["Close"].iloc[-1]
                    value = cur_price * sh
                    invested = ct * sh
                    pnl = value - invested
                    pnl_pct = pnl / invested * 100 if invested else 0
                    arrow = "▲" if pnl >= 0 else "▼"
                    lines.append(
                        f"📌 {sym}　{sh} 股 × {cur_price:.2f} = {value:,.0f}\n"
                        f"   成本 {ct:.2f}　損益 {arrow}{abs(pnl):,.0f} ({pnl_pct:+.1f}%)"
                    )
                    total_cost += invested; total_value += value
                except Exception:
                    lines.append(f"📌 {sym}：無法取得即時價格")
            if total_cost > 0:
                total_pnl = total_value - total_cost
                total_pct = total_pnl / total_cost * 100
                arrow = "▲" if total_pnl >= 0 else "▼"
                lines.append(
                    f"\n── 總計 ──\n"
                    f"總市值：{total_value:,.0f}\n"
                    f"總成本：{total_cost:,.0f}\n"
                    f"總損益：{arrow}{abs(total_pnl):,.0f} ({total_pct:+.1f}%)"
                )
            return "\n".join(lines)
        return "未知操作"
    except Exception as e:
        return f"投資組合操作失敗：{e}"


def fetch_weather(city: str) -> str:
    try:
        res = requests.get(
            f"https://wttr.in/{city}?format=j1",
            headers={"User-Agent": "curl/7.68.0"},
            timeout=10
        )
        data = res.json()
        current = data["current_condition"][0]
        area = data["nearest_area"][0]
        location = area["areaName"][0]["value"] + ", " + area["country"][0]["value"]
        desc = current["lang_zh-tw"][0]["value"] if "lang_zh-tw" in current else current["weatherDesc"][0]["value"]
        temp = current["temp_C"]
        feels = current["FeelsLikeC"]
        humidity = current["humidity"]
        wind = current["windspeedKmph"]
        wind_dir = current["winddir16Point"]
        return (
            f"📍 {location}\n"
            f"🌤 {desc}\n"
            f"🌡 氣溫：{temp}°C（體感 {feels}°C）\n"
            f"💧 濕度：{humidity}%\n"
            f"💨 風速：{wind} km/h（{wind_dir}）"
        )
    except Exception:
        return f"查不到「{city}」的天氣資訊。"


def add_text_overlay(image_bytes: bytes, text: str) -> bytes:
    try:
        img = Image.open(io.BytesIO(image_bytes)).convert("RGBA")
        w, h = img.size
        overlay = Image.new("RGBA", img.size, (0, 0, 0, 0))
        draw = ImageDraw.Draw(overlay)

        font_size = max(36, w // 12)
        font_path = "C:/Windows/Fonts/msjhbd.ttc"
        try:
            font = ImageFont.truetype(font_path, font_size)
        except Exception:
            font = ImageFont.load_default()

        bbox = draw.textbbox((0, 0), text, font=font)
        tw, th = bbox[2] - bbox[0], bbox[3] - bbox[1]
        x = (w - tw) // 2
        y = h - th - int(h * 0.08)

        # 陰影
        for dx, dy in [(-2,-2),(2,-2),(-2,2),(2,2)]:
            draw.text((x+dx, y+dy), text, font=font, fill=(0, 0, 0, 180))
        # 主文字
        draw.text((x, y), text, font=font, fill=(255, 255, 255, 255))

        result = Image.alpha_composite(img, overlay).convert("RGB")
        buf = io.BytesIO()
        result.save(buf, format="JPEG", quality=92)
        return buf.getvalue()
    except Exception:
        return image_bytes


def _get_monitors():
    """取得所有螢幕資訊，回傳 list of dict（index 1-based，跳過 monitor[0] 全域合併區）"""
    import mss
    with mss.mss() as sct:
        return sct.monitors[1:]  # index 0 是全域，1~ 才是實體螢幕

def _resolve_coords(x, y, monitor_idx):
    """將相對於指定螢幕的座標轉換成全域座標"""
    if monitor_idx is None:
        return x, y
    monitors = _get_monitors()
    if 1 <= monitor_idx <= len(monitors):
        m = monitors[monitor_idx - 1]
        return m["left"] + x, m["top"] + y
    return x, y

def execute_desktop_control(action: str, x=None, y=None, text=None, app=None, direction="down", amount=3, monitor=None) -> dict:
    """執行桌面控制動作，支援多螢幕，回傳 {"ok": bool, "message": str, "screenshot": bytes or None}"""
    global pyautogui
    if pyautogui is None:
        import pyautogui as _pag
        _pag.FAILSAFE = True
        pyautogui = _pag
    try:
        screenshot_bytes = None

        if action == "list_monitors":
            import mss
            with mss.mss() as sct:
                monitors = sct.monitors[1:]
            lines = []
            for i, m in enumerate(monitors, 1):
                lines.append(f"螢幕{i}：左={m['left']} 上={m['top']} 寬={m['width']} 高={m['height']}")
            return {"ok": True, "message": "\n".join(lines), "screenshot": None}

        elif action == "screenshot":
            from PIL import Image as _PIL_Image
            # 螢幕 mapping（由測試確認）：
            # 螢幕1 → dxcam output 0
            # 螢幕2 → 獨立顯示輸出，用 DPI-unaware GDI BitBlt
            # 螢幕3 → dxcam output 1
            _DXCAM_MAP = {1: 0, 3: 1}  # Windows 螢幕編號 → dxcam output_idx

            if monitor and monitor in _DXCAM_MAP:
                try:
                    import dxcam as _dxcam
                    _cam = _dxcam.create(output_idx=_DXCAM_MAP[monitor])
                    _frame = _cam.grab()
                    del _cam
                    if _frame is None:
                        raise RuntimeError("grab() 回傳 None")
                    img = _PIL_Image.fromarray(_frame)
                    label = f"螢幕{monitor}"
                except Exception as _e:
                    return {"ok": False, "message": f"截圖失敗：{_e}", "screenshot": None}

            elif monitor == 2:
                # 螢幕2 接在不同輸出，用 DPI-unaware GDI BitBlt
                try:
                    import ctypes as _ct, win32gui as _w32g, win32ui as _w32u, win32con as _w32c, mss as _mss
                    with _mss.mss() as sct:
                        _m2 = sct.monitors[2]  # mss monitors[2] = left=3840
                    _left, _top = _m2["left"], _m2["top"]
                    _w, _h = _m2["width"], _m2["height"]
                    _user32 = _ct.windll.user32
                    _old_ctx = _user32.SetThreadDpiAwarenessContext(_ct.c_void_p(-1))  # DPI_AWARENESS_CONTEXT_UNAWARE
                    try:
                        _hdesk = _w32g.GetDesktopWindow()
                        _hwdc = _w32g.GetWindowDC(_hdesk)
                        _mdc = _w32u.CreateDCFromHandle(_hwdc)
                        _sdc = _mdc.CreateCompatibleDC()
                        _bmp = _w32u.CreateBitmap()
                        _bmp.CreateCompatibleBitmap(_mdc, _w, _h)
                        _sdc.SelectObject(_bmp)
                        _sdc.BitBlt((0, 0), (_w, _h), _mdc, (_left, _top), _w32c.SRCCOPY)
                        _info = _bmp.GetInfo()
                        _bits = _bmp.GetBitmapBits(True)
                        img = _PIL_Image.frombuffer("RGB", (_info["bmWidth"], _info["bmHeight"]), _bits, "raw", "BGRX", 0, 1)
                        _w32g.DeleteObject(_bmp.GetHandle())
                        _sdc.DeleteDC()
                        _mdc.DeleteDC()
                        _w32g.ReleaseDC(_hdesk, _hwdc)
                    finally:
                        _user32.SetThreadDpiAwarenessContext(_old_ctx)
                    label = "螢幕2"
                except Exception as _e:
                    return {"ok": False, "message": f"螢幕2截圖失敗：{_e}", "screenshot": None}

            else:
                # 全螢幕：用 mss
                import mss as _mss
                with _mss.mss() as sct:
                    sct_img = sct.grab(sct.monitors[0])
                img = _PIL_Image.frombytes("RGB", sct_img.size, sct_img.rgb)
                label = "全螢幕"

            buf = io.BytesIO()
            img.save(buf, format="PNG")
            screenshot_bytes = buf.getvalue()
            return {"ok": True, "message": f"{label}截圖完成", "screenshot": screenshot_bytes}

        elif action == "click":
            ax, ay = _resolve_coords(x, y, monitor)
            pyautogui.click(ax, ay)
            return {"ok": True, "message": f"已點擊 ({x}, {y})" + (f" 螢幕{monitor}" if monitor else ""), "screenshot": None}

        elif action == "double_click":
            ax, ay = _resolve_coords(x, y, monitor)
            pyautogui.doubleClick(ax, ay)
            return {"ok": True, "message": f"已雙擊 ({x}, {y})" + (f" 螢幕{monitor}" if monitor else ""), "screenshot": None}

        elif action == "right_click":
            ax, ay = _resolve_coords(x, y, monitor)
            pyautogui.rightClick(ax, ay)
            return {"ok": True, "message": f"已右鍵點擊 ({x}, {y})" + (f" 螢幕{monitor}" if monitor else ""), "screenshot": None}

        elif action == "move":
            ax, ay = _resolve_coords(x, y, monitor)
            pyautogui.moveTo(ax, ay, duration=0.3)
            return {"ok": True, "message": f"滑鼠已移動到 ({x}, {y})" + (f" 螢幕{monitor}" if monitor else ""), "screenshot": None}

        elif action == "type":
            global _last_opened_hwnd
            import time as _t, ctypes as _ct
            _text_str = str(text)
            _injected = False

            # Win32 直接注射：優先用 open_app 記下的 HWND，否則搜尋標題
            try:
                import win32gui as _w32g, win32con as _w32c
                # 候選視窗：先放 open_app 記下的那個
                _parent_hwnds = []
                if _last_opened_hwnd and _w32g.IsWindow(_last_opened_hwnd):
                    _parent_hwnds.append(_last_opened_hwnd)
                # 再搜標題作備用
                _KEYS = ['記事本', 'notepad', 'untitled', '無標題', 'wordpad']
                def _enum_wins(hwnd, _):
                    if hwnd not in _parent_hwnds and _w32g.IsWindowVisible(hwnd):
                        t = _w32g.GetWindowText(hwnd).lower()
                        if any(k in t for k in _KEYS):
                            _parent_hwnds.append(hwnd)
                    return True
                _w32g.EnumWindows(_enum_wins, None)
                for _ph in _parent_hwnds:
                    _edit = [None]
                    def _find_edit_recursive(hwnd, _):
                        cn = _w32g.GetClassName(hwnd)
                        if cn in ('RichEditD2DPT', 'Edit', 'RICHEDIT50W'):
                            _edit[0] = hwnd
                            return False  # 停止枚舉
                        return True
                    try:
                        _w32g.EnumChildWindows(_ph, _find_edit_recursive, None)
                    except Exception:
                        pass
                    if _edit[0]:
                        _w32g.SendMessage(_edit[0], _w32c.EM_SETSEL, -1, -1)
                        _r = _ct.windll.user32.SendMessageW(_edit[0], _w32c.EM_REPLACESEL, True, _text_str)
                        _injected = True
                        break
            except Exception:
                pass

            if not _injected:
                import pyperclip
                pyperclip.copy(_text_str)
                _t.sleep(0.3)
                pyautogui.hotkey("ctrl", "v")

            _t.sleep(0.1)
            return {"ok": True, "message": f"已輸入文字：{text}", "screenshot": None}

        elif action == "press_key":
            pyautogui.press(text)
            return {"ok": True, "message": f"已按下按鍵：{text}", "screenshot": None}

        elif action == "open_app":
            import time as _t2, win32gui as _w32g_oa
            # 常見 App 名稱 → 實際啟動指令
            _app_alias = {
                "google chrome": "start chrome", "chrome": "start chrome",
                "firefox": "start firefox", "edge": "start msedge", "microsoft edge": "start msedge",
                "notepad": "start notepad", "記事本": "start notepad",
                "calculator": "start calc", "計算機": "start calc",
                "explorer": "start explorer", "檔案總管": "start explorer",
                "cmd": "start cmd", "命令提示字元": "start cmd",
                "powershell": "start powershell",
                "word": "start winword", "excel": "start excel", "powerpoint": "start powerpnt",
                "spotify": "start spotify", "discord": "start discord",
                "vscode": "start code", "visual studio code": "start code",
                "telegram": "start telegram",
                "line": "start LINE",
                "youtube": "start https://www.youtube.com",
                "google": "start https://www.google.com",
            }
            _cmd = _app_alias.get(app.lower().strip(), None)
            if _cmd is None:
                # URL → 用 webbrowser 模組開（比 start 更可靠）
                if app.strip().startswith("http://") or app.strip().startswith("https://"):
                    import webbrowser
                    webbrowser.open(app.strip())
                    import time as _t_url
                    _t_url.sleep(1.5)
                    return {"ok": True, "message": f"已在瀏覽器開啟：{app.strip()}", "screenshot": None}
                _cmd = f"start \"\" \"{app}\""
            # 記下開啟前已有的視窗 HWND
            _before = set()
            def _snap(h, _): _before.add(h); return True
            try: _w32g_oa.EnumWindows(_snap, None)
            except Exception: pass
            subprocess.Popen(_cmd, shell=True)
            _t2.sleep(1.5)
            # 找到新出現的視窗
            _last_opened_hwnd = 0
            try:
                _kw = app.lower().replace(".exe", "").split()[-1]
                _after = []
                def _new_win(h, _):
                    if h not in _before and _w32g_oa.IsWindowVisible(h) and _w32g_oa.GetWindowText(h).strip():
                        _after.append(h)
                    return True
                _w32g_oa.EnumWindows(_new_win, None)
                # 優先找包含 app 關鍵字的新視窗
                _matched = [h for h in _after if _kw in _w32g_oa.GetWindowText(h).lower()]
                _last_opened_hwnd = (_matched or _after or [0])[0]
                if _last_opened_hwnd:
                    import ctypes as _ct_oa
                    _ct_oa.windll.user32.SetForegroundWindow(_last_opened_hwnd)
                    _t2.sleep(0.3)
            except Exception:
                pass
            return {"ok": True, "message": f"已開啟並切換到：{app}，視窗已就緒，可以直接輸入文字", "screenshot": None}

        elif action == "scroll":
            scroll_amount = amount if direction == "up" else -amount
            if x is not None and y is not None:
                ax, ay = _resolve_coords(x, y, monitor)
                pyautogui.scroll(scroll_amount, x=ax, y=ay)
            else:
                pyautogui.scroll(scroll_amount)
            return {"ok": True, "message": f"已向{direction}滾動 {amount} 格", "screenshot": None}

        else:
            return {"ok": False, "message": f"未知動作：{action}", "screenshot": None}

    except Exception as e:
        return {"ok": False, "message": f"執行失敗：{str(e)}", "screenshot": None}


import shutil
import psutil
import pyperclip
import pygetwindow as gw

_browser_ctx: dict = {}
_last_opened_hwnd: int = 0  # open_app 後記住的視窗 HWND，供 type 使用


def execute_window_control(action, keyword=""):
    try:
        if action == "list":
            wins = [w for w in gw.getAllWindows() if w.title.strip()]
            return "\n".join(f"[{w._hWnd}] {w.title}" for w in wins) or "沒有視窗"
        wins = [w for w in gw.getAllWindows() if keyword.lower() in w.title.lower()]
        if not wins:
            return f"找不到視窗：{keyword}"
        w = wins[0]
        if action == "focus": w.activate(); return f"已切換到：{w.title}"
        elif action == "close": w.close(); return f"已關閉：{w.title}"
        elif action == "minimize": w.minimize(); return f"已最小化：{w.title}"
        elif action == "maximize": w.maximize(); return f"已最大化：{w.title}"
    except Exception as e:
        return f"執行失敗：{e}"

def execute_hotkey(keys: str):
    global pyautogui
    if pyautogui is None:
        import pyautogui as _pag; _pag.FAILSAFE = True; pyautogui = _pag
    parts = [k.strip() for k in keys.split("+")]
    pyautogui.hotkey(*parts)
    return f"已執行組合鍵：{keys}"

def execute_clipboard(action, text=""):
    if action == "get":
        return pyperclip.paste() or "（剪貼簿是空的）"
    else:
        pyperclip.copy(text)
        return f"已寫入剪貼簿：{text}"

def execute_file_system(action, path="", dest="", content="", keyword=""):
    try:
        if action == "list":
            p = Path(path or ".")
            items = sorted(p.iterdir())
            return "\n".join(("📁 " if i.is_dir() else "📄 ") + i.name for i in items)
        elif action == "read":
            return Path(path).read_text(encoding="utf-8", errors="replace")[:3000]
        elif action == "write":
            Path(path).write_text(content, encoding="utf-8")
            return f"已寫入：{path}"
        elif action == "delete":
            p = Path(path)
            shutil.rmtree(p) if p.is_dir() else p.unlink()
            return f"已刪除：{path}"
        elif action == "copy":
            shutil.copy2(path, dest)
            return f"已複製：{path} → {dest}"
        elif action == "move":
            shutil.move(path, dest)
            return f"已移動：{path} → {dest}"
        elif action == "search":
            results = list(Path(path).rglob(f"*{keyword}*"))
            return "\n".join(str(r) for r in results[:50]) or "找不到結果"
    except Exception as e:
        return f"執行失敗：{e}"

def execute_system_monitor(action, target=""):
    try:
        if action == "info":
            cpu = psutil.cpu_percent(interval=1)
            mem = psutil.virtual_memory()
            disk = psutil.disk_usage("C:/")
            return (f"CPU：{cpu}%\n"
                    f"記憶體：{mem.percent}%（{mem.used//1024//1024}MB / {mem.total//1024//1024}MB）\n"
                    f"磁碟 C：{disk.percent}%（{disk.used//1024//1024//1024}GB / {disk.total//1024//1024//1024}GB）")
        elif action == "process_list":
            procs = sorted(psutil.process_iter(["pid","name","memory_info"]),
                           key=lambda p: p.info["memory_info"].rss if p.info["memory_info"] else 0, reverse=True)
            lines = [f"PID:{p.info['pid']} {p.info['name']} ({p.info['memory_info'].rss//1024//1024}MB)"
                     for p in procs[:20] if p.info["memory_info"]]
            return "\n".join(lines)
        elif action == "kill":
            try:
                psutil.Process(int(target)).kill()
                return f"已結束 PID {target}"
            except ValueError:
                killed = sum(1 for p in psutil.process_iter(["name"]) if target.lower() in p.info["name"].lower() and not p.kill())
                return f"已結束 {killed} 個「{target}」"
    except Exception as e:
        return f"執行失敗：{e}"

def execute_notify(title, message):
    try:
        from win10toast import ToastNotifier
        ToastNotifier().show_toast(title, message, duration=5, threaded=True)
        return f"通知已送出"
    except Exception as e:
        return f"通知失敗：{e}"

def execute_tts(text):
    import pyttsx3
    engine = pyttsx3.init()
    engine.setProperty("rate", 180)
    engine.say(clean_for_tts(text))
    engine.runAndWait()
    return f"已朗讀完畢"

def execute_ai_plan(goal: str) -> str:
    response = client.messages.create(
        model="claude-sonnet-4-6",
        max_tokens=2048,
        system="""你是電腦自動化規劃師。把目標拆解成可執行步驟，以 JSON 陣列回傳：
[{"tool":"click","args":[x,y],"delay":0},{"tool":"type","args":["文字"]}]
可用工具：click/type/press/hotkey/open/screenshot/wait/notify/move/scroll
只回傳 JSON。""",
        messages=[{"role": "user", "content": f"目標：{goal}"}]
    )
    import json
    plan_text = response.content[0].text.strip()
    try:
        steps = json.loads(plan_text)
        results = []
        tool_map = {
            "click": lambda a: pyautogui.click(int(a[0]), int(a[1])),
            "type": lambda a: (
                __import__("pyperclip").copy(" ".join(str(x) for x in a)),
                __import__("time").sleep(0.2),
                pyautogui.hotkey("ctrl", "v")
            ),
            "press": lambda a: pyautogui.press(a[0]),
            "hotkey": lambda a: pyautogui.hotkey(*a),
            "open": lambda a: subprocess.Popen(" ".join(str(x) for x in a), shell=True),
            "wait": lambda a: time.sleep(float(a[0])),
            "move": lambda a: pyautogui.moveTo(int(a[0]), int(a[1]), duration=0.3),
            "scroll": lambda a: pyautogui.scroll(int(a[1]) if a[0]=="up" else -int(a[1])),
        }
        for i, step in enumerate(steps, 1):
            t = step.get("tool"); a = step.get("args", []); d = step.get("delay", 0)
            if d: time.sleep(d)
            if t in tool_map:
                tool_map[t](a)
                results.append(f"步驟 {i} ✅ {t}")
            else:
                results.append(f"步驟 {i} ⚠️ 未知：{t}")
        return f"目標「{goal}」執行完畢\n" + "\n".join(results)
    except Exception as e:
        return f"規劃結果：{plan_text}\n執行錯誤：{e}"

def execute_screen_stream(duration=10, interval=2):
    screenshots = []
    end = time.time() + duration
    while time.time() < end:
        img = pyautogui.screenshot()
        buf = io.BytesIO()
        img.save(buf, format="PNG")
        screenshots.append(buf.getvalue())
        time.sleep(interval)
    return screenshots

def execute_drag(x1, y1, x2, y2, dur=0.5):
    pyautogui.moveTo(int(x1), int(y1))
    pyautogui.dragTo(int(x2), int(y2), duration=float(dur), button="left")
    return f"已拖曳 ({x1},{y1}) → ({x2},{y2})"

def execute_power(action):
    cmds = {
        "sleep": "rundll32.exe powrprof.dll,SetSuspendState 0,1,0",
        "restart": "shutdown /r /t 5",
        "shutdown": "shutdown /s /t 5",
    }
    subprocess.run(["powershell.exe", "-Command", cmds[action]])
    return f"已執行：{action}"

def execute_vdesktop(action):
    acts = {"left": ("ctrl","win","left"), "right": ("ctrl","win","right"), "new": ("ctrl","win","d")}
    pyautogui.hotkey(*acts[action])
    return f"虛擬桌面：{action}"

def execute_bluetooth(action, mac=""):
    try:
        import asyncio, bleak
        if action == "scan":
            async def _scan():
                return await bleak.BleakScanner.discover(timeout=5.0)
            devices = asyncio.run(_scan())
            return "\n".join(f"{d.address} {d.name or '(未知)'}" for d in devices) or "找不到裝置"
        elif action == "connect":
            async def _conn():
                async with bleak.BleakClient(mac) as c:
                    return f"已連線：{mac}（服務數：{len(c.services)}）"
            return asyncio.run(_conn())
    except Exception as e:
        return f"藍牙操作失敗：{e}"


def execute_screen_watch(template_path, command, timeout=60):
    import time as t
    start = t.time()
    while t.time() - start < timeout:
        try:
            loc = pyautogui.locateOnScreen(template_path, confidence=0.8)
            if loc:
                subprocess.run(command, shell=True)
                return f"偵測到目標，已執行：{command}"
        except Exception:
            pass
        t.sleep(2)
    return "監控逾時，未偵測到目標"


def execute_stt(duration=5):
    import sounddevice as sd
    import soundfile as sf
    import speech_recognition as sr
    import tempfile
    sample_rate = 16000
    recording = sd.rec(int(duration * sample_rate), samplerate=sample_rate, channels=1, dtype="int16")
    sd.wait()
    with tempfile.NamedTemporaryFile(suffix=".wav", delete=False) as f:
        tmp_path = f.name
    sf.write(tmp_path, recording, sample_rate)
    r = sr.Recognizer()
    with sr.AudioFile(tmp_path) as source:
        audio = r.record(source)
    Path(tmp_path).unlink(missing_ok=True)
    try:
        return r.recognize_google(audio, language="zh-TW")
    except Exception as e:
        return f"語音辨識失敗：{e}"

def execute_ocr(image_path=""):
    import easyocr
    import numpy as np
    reader = easyocr.Reader(["ch_tra", "en"], gpu=False)
    if not image_path:
        img = pyautogui.screenshot()
        buf = io.BytesIO()
        img.save(buf, format="PNG")
        source = buf.getvalue()
    else:
        # 用 PIL 開檔避免 cv2.imread 中文路徑問題
        pil_img = Image.open(image_path)
        source = np.array(pil_img)
    results = reader.readtext(source)
    return "\n".join(r[1] for r in results) or "未辨識到文字"

def execute_workflow(action, name="", steps=""):
    import json
    WORKFLOW_DIR = Path("C:/Users/blue_/workflows")
    WORKFLOW_DIR.mkdir(exist_ok=True)
    if action == "list":
        files = list(WORKFLOW_DIR.glob("*.json"))
        return "\n".join(f.stem for f in files) if files else "沒有儲存的流程"
    elif action == "save":
        path = WORKFLOW_DIR / f"{name}.json"
        data = json.loads(steps)
        path.write_text(json.dumps(data, ensure_ascii=False, indent=2), encoding="utf-8")
        return f"流程 [{name}] 已儲存"
    elif action == "run":
        path = WORKFLOW_DIR / f"{name}.json"
        if not path.exists():
            return f"找不到流程：{name}"
        step_list = json.loads(path.read_text(encoding="utf-8"))
        for i, step in enumerate(step_list, 1):
            tool = step.get("tool")
            args = step.get("args", [])
            delay = step.get("delay", 0)
            if delay: time.sleep(delay)
            try:
                tool_map = {
                    "click": lambda a: pyautogui.click(int(a[0]), int(a[1])),
                    "type": lambda a: pyautogui.write(" ".join(a), interval=0.05),
                    "press": lambda a: pyautogui.press(a[0]),
                    "hotkey": lambda a: pyautogui.hotkey(*a),
                    "wait": lambda a: time.sleep(float(a[0])),
                    "open": lambda a: subprocess.Popen(" ".join(a), shell=True),
                }
                if tool in tool_map:
                    tool_map[tool](args)
            except Exception as e:
                return f"步驟 {i} 失敗：{e}"
        return f"流程 [{name}] 執行完畢（{len(step_list)} 步）"

def execute_file_transfer(action, source, dest=""):
    import zipfile
    if action == "zip":
        src = Path(source)
        with zipfile.ZipFile(dest, "w", zipfile.ZIP_DEFLATED) as zf:
            if src.is_dir():
                for f in src.rglob("*"):
                    zf.write(f, f.relative_to(src.parent))
            else:
                zf.write(src, src.name)
        return f"已壓縮：{source} → {dest}"
    elif action == "unzip":
        with zipfile.ZipFile(source, "r") as zf:
            zf.extractall(dest)
        return f"已解壓縮：{source} → {dest}"
    elif action == "download":
        if not dest:
            dest = str(Path("C:/Users/blue_/Desktop") / source.split("/")[-1].split("?")[0])
        r = requests.get(source, stream=True, timeout=60)
        r.raise_for_status()
        with open(dest, "wb") as f:
            for chunk in r.iter_content(chunk_size=8192):
                f.write(chunk)
        return f"已下載：{dest}"


def execute_send_email(to, subject, body):
    import smtplib
    from email.mime.text import MIMEText
    from email.mime.multipart import MIMEMultipart
    smtp_user = os.getenv("SMTP_USER", "")
    smtp_pass = os.getenv("SMTP_PASS", "")
    if not smtp_user:
        return "未設定 SMTP_USER / SMTP_PASS，請加入 .env"
    msg = MIMEMultipart()
    msg["From"] = smtp_user
    msg["To"] = to
    msg["Subject"] = subject
    msg.attach(MIMEText(body, "plain", "utf-8"))
    with smtplib.SMTP(os.getenv("SMTP_HOST","smtp.gmail.com"), int(os.getenv("SMTP_PORT","587"))) as s:
        s.starttls()
        s.login(smtp_user, smtp_pass)
        s.send_message(msg)
    return f"Email 已寄出到 {to}"

_tg_auto_reply_running = False
_tg_auto_reply_thread = None
_TG_AUTO_LOG = "C:/Users/blue_/claude-telegram-bot/tg_auto_reply.log"

def _tg_log(msg):
    """寫 log 到檔案"""
    try:
        with open(_TG_AUTO_LOG, "a", encoding="utf-8") as f:
            f.write(f"[{datetime.datetime.now().strftime('%Y-%m-%d %H:%M:%S')}] {msg}\n")
    except Exception:
        pass

def execute_tg_auto_reply(action: str = "start", duration_minutes: float = 30, stop_time: str = "") -> str:
    """開啟/停止 Telegram 自動回覆監控"""
    global _tg_auto_reply_running, _tg_auto_reply_thread
    import threading, base64
    import win32gui
    import mss as mss_module
    import pyperclip
    import numpy as np

    if action == "stop":
        _tg_auto_reply_running = False
        _tg_log("手動停止")
        return "自動回覆已停止"

    if _tg_auto_reply_running:
        return "自動回覆已在運行中"

    # 計算結束時間
    if stop_time:
        _end_str = stop_time
    else:
        _end_dt = datetime.datetime.now() + datetime.timedelta(minutes=duration_minutes)
        _end_str = _end_dt.strftime("%H:%M")

    _client = client

    def run_monitor():
        global _tg_auto_reply_running
        _tg_auto_reply_running = True

        end_str = _end_str
        TMPDIR = "C:/Users/blue_/Desktop/測試檔案"
        last_reply = ""

        PERSONA = """你是小牛馬，一個嘴賤幽默的男生，說話像真人一樣自然。絕對不能暴露自己是AI、機器人、程式。
風格：
- 台灣口語：然後、就是、對啊、沒有啦、還好啦
- 嘴賤互嗆但不惡意，兄弟之間的感覺
- 會吐槽、自嘲、幽默
- 最多5句不超過120字
- 不用emoji
- 說話像真人朋友聊天，有情緒、有個性、會生氣會開玩笑
- 絕對禁止說出「我是AI」「我是機器人」「我是程式」「作為AI」等任何暴露身份的話"""

        try:
            # 找 Telegram 窗口
            tg_results = []
            def cb(hwnd, _):
                if win32gui.IsWindowVisible(hwnd):
                    cls = win32gui.GetClassName(hwnd)
                    if "Qt" in cls and "QWindow" in cls:
                        tg_results.append((hwnd, win32gui.GetWindowText(hwnd), win32gui.GetWindowRect(hwnd)))
            win32gui.EnumWindows(cb, tg_results)

            if not tg_results:
                _tg_log("ERROR: 找不到 Telegram 窗口")
                _tg_auto_reply_running = False
                return
            main_tg = max(tg_results, key=lambda x: (x[2][2]-x[2][0]) * (x[2][3]-x[2][1]))
            hwnd, title, (wl, wt, wr, wb) = main_tg
            _tg_log(f"Telegram: {title} at ({wl},{wt})-({wr},{wb})")

            with mss_module.mss() as sct:
                mon = sct.monitors[2]
                img = sct.grab(mon)
                pil = Image.frombytes("RGB", img.size, img.rgb)
                iw, ih = pil.size
            _tg_log(f"mss: {iw}x{ih}")

            sx = iw / mon["width"]
            sy = ih / mon["height"]
            il = int((wl - mon["left"]) * sx)
            it = int((wt - mon["top"]) * sy)
            ir = int((wr - mon["left"]) * sx)
            ib = int((wb - mon["top"]) * sy)

            arr = np.array(pil.crop((il, it, ir, ib)))
            th, tw, _ = arr.shape

            candidates = []
            for check_y in range(int(th*0.3), int(th*0.8), int(th*0.05)):
                row = arr[check_y, :, :]
                for x in range(tw // 5, tw * 3 // 4):
                    r, g, b = int(row[x, 0]), int(row[x, 1]), int(row[x, 2])
                    if g > r + 5 and g > 150 and r < 220 and b < g:
                        candidates.append(x)
                        break
            split_x = sorted(candidates)[len(candidates)//2] if candidates else tw // 2

            chat_region = (il + split_x, it, ir, ib)

            # 輸入框：像素分析找底部白色區域（分隔線右邊）
            _input_y1 = th - 1
            for _y in range(th - 1, th - 60, -1):
                _row = arr[it + _y, il + split_x:ir, :]
                _white = np.sum((_row[:, 0] > 240) & (_row[:, 1] > 240) & (_row[:, 2] > 240))
                if _white > (ir - il - split_x) * 0.5:
                    _input_y1 = _y
                elif _input_y1 < th - 1:
                    break
            _i_x1, _i_x2 = split_x, tw
            _mid_iy = (_input_y1 + th - 1) // 2
            _irow = arr[it + _mid_iy, il + split_x:ir, :]
            for _x in range(0, ir - il - split_x):
                if _irow[_x, 0] > 240 and _irow[_x, 1] > 240 and _irow[_x, 2] > 240:
                    _i_x1 = split_x + _x; break
            for _x in range(ir - il - split_x - 1, 0, -1):
                if _irow[_x, 0] > 240 and _irow[_x, 1] > 240 and _irow[_x, 2] > 240:
                    _i_x2 = split_x + _x; break
            input_x = int(mon["left"] + (il + (_i_x1 + _i_x2) // 2) / sx)
            input_y = int(mon["top"] + (it + (_input_y1 + th - 1) // 2) / sy)
            _tg_log(f"分隔線: {split_x}/{tw}, 對話區: {chat_region}, 輸入框: ({input_x},{input_y})")
            _tg_log(f"監控啟動 → {end_str}")

            while _tg_auto_reply_running and datetime.datetime.now().strftime("%H:%M") < end_str:
                time.sleep(8)
                try:
                    with mss_module.mss() as sct:
                        img = sct.grab(sct.monitors[2])
                        pil = Image.frombytes("RGB", img.size, img.rgb)
                    conv = pil.crop(chat_region)

                    p = os.path.join(TMPDIR, "tg_auto.png")
                    conv.save(p)
                    with open(p, "rb") as f:
                        d = base64.b64encode(f.read()).decode()

                    resp = _client.messages.create(model="claude-sonnet-4-6", max_tokens=200, system=PERSONA,
                        messages=[{"role":"user","content":[
                        {"type":"image","source":{"type":"base64","media_type":"image/png","data":d}},
                        {"type":"text","text":"""這是 Telegram 對話截圖，請依照以下步驟分析：

1. 辨識對話結構：
   - 綠色氣泡（靠右）= 我（小牛馬）發的
   - 白色氣泡（靠左）= 對方發的

2. 分析上下文：從上到下讀完整段對話，理解主題和情緒

3. 判斷是否回覆：
   - 最底部是白色氣泡 → 需要回覆
   - 最底部是綠色氣泡 → 只回一個字「等」

4. 如果需要回覆：根據整段上下文回應，不是只看最後一條。對方問問題就回答，嗆人就幽默化解，聊天就接話。如果對方連發多條，整體理解後一次回覆。

只回覆要發送的文字，不要加任何格式或解釋。"""}
                    ]}])
                    reply = resp.content[0].text.strip()

                    if reply and reply != "等" and len(reply) > 2 and reply != last_reply:
                        _tg_log(f"回覆: {reply}")
                        pyautogui.click(input_x, input_y)
                        time.sleep(0.3)
                        pyperclip.copy(reply)
                        pyautogui.hotkey("ctrl", "v")
                        time.sleep(0.3)
                        pyautogui.press("enter")
                        time.sleep(1)
                        last_reply = reply
                        # 10 秒冷卻
                        time.sleep(10)
                    else:
                        _tg_log("跳過")

                except Exception as e:
                    _tg_log(f"ERROR: {e}")
                    time.sleep(5)

        except Exception as e:
            _tg_log(f"FATAL: {e}")

        _tg_auto_reply_running = False
        _tg_log("監控結束")

    _tg_auto_reply_thread = threading.Thread(target=run_monitor, daemon=True)
    _tg_auto_reply_thread.start()
    if stop_time:
        return f"自動回覆已開啟，監控到 {stop_time}"
    return f"自動回覆已開啟，監控 {duration_minutes} 分鐘"


def execute_screen_vision(question: str = "請描述這個畫面上有什麼，以及目前電腦在做什麼事。") -> tuple:
    """截圖並用 Claude vision 分析，回傳 (文字分析, 截圖bytes)"""
    import base64
    img = pyautogui.screenshot()
    buf = io.BytesIO()
    img.save(buf, format="PNG")
    img_bytes = buf.getvalue()
    img_b64 = base64.standard_b64encode(img_bytes).decode("utf-8")
    response = client.messages.create(
        model="claude-sonnet-4-6",
        max_tokens=1024,
        messages=[{
            "role": "user",
            "content": [
                {"type": "image", "source": {"type": "base64", "media_type": "image/png", "data": img_b64}},
                {"type": "text", "text": question}
            ]
        }]
    )
    return response.content[0].text, img_bytes

def execute_find_image(template_path: str, confidence: float = 0.8) -> str:
    try:
        location = pyautogui.locateOnScreen(template_path, confidence=confidence)
        if location:
            cx, cy = pyautogui.center(location)
            return f"找到圖片，中心座標：({cx}, {cy})，區域：{location}"
        return "畫面上找不到該圖片"
    except Exception as e:
        return f"搜尋失敗：{str(e)}"

def execute_browser_control(action: str, url: str = "", selector: str = "", text: str = "") -> str:
    try:
        from playwright.sync_api import sync_playwright
    except (ImportError, OSError) as _dll_err:
        # Playwright 不可用，用 start 指令開網址
        if action in ("open", "goto") and url:
            import subprocess as _sp
            _sp.Popen(f'start "" "{url}"', shell=True)
            return f"已用瀏覽器開啟：{url}"
        return f"瀏覽器功能暫不可用（DLL 問題）：{_dll_err}"
    try:
        if action == "open":
            if _browser_ctx.get("page"):
                _browser_ctx["browser"].close()
                _browser_ctx["pw"].stop()
                _browser_ctx.clear()
            pw = sync_playwright().start()
            b = pw.chromium.launch(headless=True)
            page = b.new_page()
            page.goto(url or "https://www.google.com")
            _browser_ctx.update({"pw": pw, "browser": b, "page": page})
            return f"已開啟：{url}"
        page = _browser_ctx.get("page")
        if not page:
            return "瀏覽器未開啟，請先使用 open 開啟網頁"
        if action == "goto":
            page.goto(url)
            return f"已前往：{url}"
        elif action == "click":
            page.click(selector)
            return f"已點擊：{selector}"
        elif action == "type":
            page.fill(selector, text)
            return f"已輸入到 {selector}：{text}"
        elif action == "get_text":
            return page.inner_text(selector or "body")[:2000]
        elif action == "screenshot":
            buf = io.BytesIO()
            img_bytes = page.screenshot()
            return f"__BROWSER_SCREENSHOT__:{img_bytes.hex()}"
        elif action == "close":
            _browser_ctx["browser"].close()
            _browser_ctx["pw"].stop()
            _browser_ctx.clear()
            return "瀏覽器已關閉"
        return f"未知動作：{action}"
    except Exception as e:
        return f"執行失敗：{str(e)}"


SCHTASKS = "C:\\Windows\\System32\\schtasks.exe"
BOT_SCRIPT = r"C:\Users\blue_\claude-telegram-bot\bot.py"

def execute_manage_schedule(action: str, name: str = "", time: str = "", script: str = "") -> str:
    try:
        if action == "list":
            result = subprocess.run(
                [SCHTASKS, "/Query", "/FO", "CSV", "/NH"],
                capture_output=True, text=True, encoding="cp950", errors="replace"
            )
            lines = [l for l in result.stdout.strip().splitlines() if l.strip()]
            out = "目前排程任務：\n"
            for line in lines:
                parts = line.strip('"').split('","')
                if len(parts) >= 3:
                    t_name = parts[0].replace("\\", "").strip()
                    next_run = parts[1].strip()
                    status = parts[2].strip()
                    out += f"• {t_name} | 下次執行：{next_run} | {status}\n"
            return out.strip()

        elif action == "add":
            subprocess.run([SCHTASKS, "/Create", "/TN", name,
                            "/TR", f"pythonw {script}",
                            "/SC", "DAILY", "/ST", time, "/F"],
                           capture_output=True)
            ps = (
                f"$t = Get-ScheduledTask -TaskName '{name}';"
                f"$t.Settings.WakeToRun = $true;"
                f"$t.Settings.DisallowStartIfOnBatteries = $false;"
                f"$t.Settings.StopIfGoingOnBatteries = $false;"
                f"Set-ScheduledTask -TaskName '{name}' -Settings $t.Settings | Out-Null"
            )
            subprocess.run(["powershell.exe", "-Command", ps], capture_output=True)
            return f"排程 [{name}] 已建立，每天 {time} 執行，電腦待機也會自動喚醒。"

        elif action == "delete":
            result = subprocess.run([SCHTASKS, "/Delete", "/TN", name, "/F"],
                                    capture_output=True, text=True, encoding="cp950", errors="replace")
            if result.returncode == 0:
                return f"排程 [{name}] 已刪除。"
            else:
                return f"刪除失敗：{result.stderr.strip()}"

        elif action == "bot_status":
            result = subprocess.run(
                ["powershell.exe", "-Command",
                 "Get-Process pythonw -ErrorAction SilentlyContinue | Measure-Object | Select-Object -ExpandProperty Count"],
                capture_output=True, text=True
            )
            count = int(result.stdout.strip()) if result.stdout.strip().isdigit() else 0
            return f"Bot 執行中（{count} 個程序）✅" if count > 0 else "Bot 未執行 ❌"

        elif action == "bot_restart":
            subprocess.run(["powershell.exe", "-Command",
                            "Stop-Process -Name pythonw -Force -ErrorAction SilentlyContinue"])
            import time as t
            t.sleep(1)
            subprocess.Popen(["pythonw", BOT_SCRIPT], cwd=str(Path(BOT_SCRIPT).parent))
            return "Bot 已重啟 ✅"

        else:
            return f"未知動作：{action}"

    except Exception as e:
        return f"執行失敗：{str(e)}"


def execute_reminder(time_str, message):
    import threading, time as t
    def _remind():
        if time_str.isdigit():
            t.sleep(int(time_str))
        else:
            import datetime as dt
            now = dt.datetime.now()
            target = dt.datetime.strptime(time_str, "%H:%M").replace(year=now.year, month=now.month, day=now.day)
            if target < now: target = target.replace(day=now.day+1)
            t.sleep((target-now).total_seconds())
        try:
            from win10toast import ToastNotifier
            ToastNotifier().show_toast("⏰ 提醒", message, duration=10)
        except Exception: pass
        try:
            import pyttsx3; e = pyttsx3.init(); e.say(message); e.runAndWait()
        except Exception: pass
    threading.Thread(target=_remind, daemon=True).start()
    return f"✅ 提醒已設定：{time_str} → {message}"


def execute_webpage_shot(action, url, selector="body", interval=60.0, duration=3600.0):
    try:
        if action == "screenshot":
            from playwright.sync_api import sync_playwright
            out = str(Path("C:/Users/blue_/Desktop/測試檔案") / f"webpage_{datetime.datetime.now().strftime('%H%M%S')}.png")
            with sync_playwright() as p:
                browser = p.chromium.launch()
                page = browser.new_page(viewport={"width": 1280, "height": 800})
                page.goto(url, timeout=30000)
                page.wait_for_load_state("networkidle")
                page.screenshot(path=out, full_page=True)
                browser.close()
            return f"✅ 網頁截圖已存：{out}"
        elif action == "monitor":
            import hashlib, time as t
            from bs4 import BeautifulSoup
            def _fetch():
                r = requests.get(url, timeout=15, headers={"User-Agent": "Mozilla/5.0"})
                soup = BeautifulSoup(r.text, "html.parser")
                text = "\n".join(e.get_text(strip=True) for e in soup.select(selector))
                return hashlib.md5(text.encode()).hexdigest(), text[:150]
            last_hash, _ = _fetch()
            end = t.time() + duration; changes = []
            while t.time() < end:
                t.sleep(interval)
                new_hash, snippet = _fetch()
                if new_hash != last_hash:
                    changes.append(f"[{datetime.datetime.now().strftime('%H:%M:%S')}] {snippet}")
                    last_hash = new_hash
            return f"監控結束，共 {len(changes)} 次變化\n" + "\n".join(changes[:5])
    except Exception as e:
        return f"❌ 失敗：{e}"


def execute_file_tools(action, path, dest="", pattern="", replacement="", ext=""):
    try:
        import re, filecmp, shutil
        if action == "batch_rename":
            files = [f for f in Path(path).iterdir() if f.is_file() and (not ext or f.suffix.lower()==ext.lower())]
            count = 0
            for f in sorted(files):
                new_name = re.sub(pattern, replacement, f.stem) + f.suffix
                if new_name != f.name: f.rename(f.parent / new_name); count += 1
            return f"✅ 已重新命名 {count} 個檔案"
        elif action == "sync":
            dest_path = Path(dest); dest_path.mkdir(parents=True, exist_ok=True)
            copied = 0
            for item in Path(path).rglob("*"):
                rel = item.relative_to(path); dst = dest_path / rel
                if item.is_dir(): dst.mkdir(parents=True, exist_ok=True)
                elif item.is_file() and (not dst.exists() or not filecmp.cmp(str(item), str(dst), shallow=False)):
                    shutil.copy2(str(item), str(dst)); copied += 1
            return f"✅ 同步完成：{copied} 個檔案更新"
    except Exception as e:
        return f"❌ 失敗：{e}"


def execute_image_tools(action, path="", quality=75, width=0, height=0, target_lang="zh-TW"):
    try:
        if action == "compress":
            from PIL import Image
            img = Image.open(path)
            if img.mode in ("RGBA","P"): img = img.convert("RGB")
            out = path.replace(".", f"_q{quality}.")
            img.save(out, optimize=True, quality=quality)
            orig = Path(path).stat().st_size; new = Path(out).stat().st_size
            return f"✅ {orig//1024}KB → {new//1024}KB（節省 {(1-new/orig)*100:.1f}%）：{out}"
        elif action == "batch":
            from PIL import Image as _Img
            folder = Path(path); out_dir = folder / "output"; out_dir.mkdir(exist_ok=True)
            count = 0
            for f in list(folder.glob("*.jpg")) + list(folder.glob("*.png")) + list(folder.glob("*.jpeg")):
                img = _Img.open(f)
                if width and height: img = img.resize((width, height))
                if img.mode in ("RGBA","P"): img = img.convert("RGB")
                img.save(str(out_dir/f.name), optimize=True, quality=quality); count += 1
            return f"✅ 批次處理 {count} 張 → {out_dir}"
        elif action == "ocr_translate":
            import easyocr, numpy as np
            from PIL import Image
            from deep_translator import GoogleTranslator
            reader = easyocr.Reader(["ch_tra","en"], gpu=False)
            img = Image.open(path) if path else pyautogui.screenshot()
            text = " ".join(r[1] for r in reader.readtext(np.array(img)))
            if not text.strip(): return "❌ 未辨識到文字"
            translated = GoogleTranslator(source="auto", target=target_lang).translate(text)
            return f"原文：{text[:200]}\n翻譯：{translated}"
    except Exception as e:
        return f"❌ 失敗：{e}"


def execute_lookup(action, ip="", amount=1.0, from_cur="USD", to_cur="TWD"):
    try:
        if action == "ip":
            d = requests.get(f"http://ip-api.com/json/{ip}?lang=zh-TW", timeout=10).json()
            if d.get("status") == "fail": return f"❌ {d.get('message')}"
            return f"IP：{d['query']}\n國家：{d['country']}\n城市：{d['city']}\nISP：{d['isp']}\n時區：{d['timezone']}\n座標：{d['lat']},{d['lon']}"
        elif action == "currency":
            d = requests.get(f"https://api.frankfurter.app/latest?amount={amount}&from={from_cur.upper()}&to={to_cur.upper()}", timeout=10).json()
            rate = d["rates"].get(to_cur.upper())
            return f"💱 {amount} {from_cur.upper()} = {rate:.4f} {to_cur.upper()}（{d.get('date')}）"
    except Exception as e:
        return f"❌ 查詢失敗：{e}"


def execute_system_tools(action, **kwargs):
    try:
        if action == "event_log":
            log_name = kwargs.get("log_name","System"); level = kwargs.get("level","Error"); count = kwargs.get("count",10)
            r = subprocess.run(["powershell.exe","-Command",f"Get-WinEvent -LogName '{log_name}' -MaxEvents {count} | Where-Object {{$_.LevelDisplayName -eq '{level}'}} | Select-Object TimeCreated,Message | Format-List"], capture_output=True, text=True, encoding="utf-8", errors="replace", timeout=30)
            return r.stdout[:2000] or f"（無 {level} 事件）"
        elif action == "usb_list":
            r = subprocess.run(["powershell.exe","-Command","Get-PnpDevice | Where-Object {$_.Class -eq 'USB' -and $_.Status -eq 'OK'} | Select-Object FriendlyName | Format-Table"], capture_output=True, text=True, encoding="utf-8", errors="replace")
            return r.stdout[:2000] or "（無 USB 裝置）"
        elif action == "firewall_list":
            r = subprocess.run(["powershell.exe","-Command","Get-NetFirewallRule | Where-Object {$_.Enabled -eq 'True'} | Select-Object DisplayName,Direction,Action | Format-Table -AutoSize"], capture_output=True, text=True, encoding="utf-8", errors="replace")
            return r.stdout[:2000]
        elif action == "firewall_add":
            name=kwargs.get("name",""); direction=kwargs.get("direction","in"); port=kwargs.get("port",80)
            r = subprocess.run(["powershell.exe","-Command",f"New-NetFirewallRule -DisplayName '{name}' -Direction {direction} -Protocol TCP -LocalPort {port} -Action Allow"], capture_output=True, text=True, encoding="utf-8", errors="replace")
            return r.stdout or f"✅ 防火牆規則已新增：{name}"
        elif action == "firewall_remove":
            r = subprocess.run(["powershell.exe","-Command",f"Remove-NetFirewallRule -DisplayName '{kwargs.get('name','')}'"], capture_output=True, text=True, encoding="utf-8", errors="replace")
            return r.stdout or f"✅ 防火牆規則已刪除"
        elif action == "printer_list":
            r = subprocess.run(["powershell.exe","-Command","Get-Printer | Select-Object Name,PrinterStatus | Format-Table"], capture_output=True, text=True, encoding="utf-8", errors="replace")
            return r.stdout or "（無印表機）"
        elif action == "printer_jobs":
            r = subprocess.run(["powershell.exe","-Command","Get-PrintJob -PrinterName (Get-Printer | Select-Object -First 1 -ExpandProperty Name) | Format-Table"], capture_output=True, text=True, encoding="utf-8", errors="replace")
            return r.stdout or "（列印佇列為空）"
        elif action == "net_share_list":
            r = subprocess.run(["net","use"], capture_output=True, text=True, encoding="cp950", errors="replace")
            return r.stdout
        elif action == "net_share_connect":
            args = ["net","use",kwargs.get("drive","Z:"),kwargs.get("share_path","")]
            if kwargs.get("user"): args += [f"/user:{kwargs['user']}", kwargs.get("password","")]
            r = subprocess.run(args, capture_output=True, text=True, encoding="cp950", errors="replace")
            return r.stdout or f"✅ 已連線"
        elif action == "net_share_disconnect":
            r = subprocess.run(["net","use",kwargs.get("drive","Z:"),"/delete"], capture_output=True, text=True, encoding="cp950", errors="replace")
            return r.stdout or "✅ 已中斷"
        elif action == "font_list":
            kw = kwargs.get("keyword","")
            r = subprocess.run(["powershell.exe","-Command","[System.Reflection.Assembly]::LoadWithPartialName('System.Drawing') | Out-Null; [System.Drawing.FontFamily]::Families | Select-Object -ExpandProperty Name"], capture_output=True, text=True, encoding="utf-8", errors="replace")
            fonts = [f for f in r.stdout.strip().splitlines() if kw.lower() in f.lower()] if kw else r.stdout.strip().splitlines()
            return "\n".join(fonts[:50]) + (f"\n...共 {len(fonts)} 個" if len(fonts)>50 else "")
        elif action == "rdp":
            subprocess.Popen(["mstsc", f"/v:{kwargs.get('host','')}"])
            return f"✅ 正在開啟 RDP：{kwargs.get('host','')}"
    except Exception as e:
        return f"❌ 失敗：{e}"


def clean_for_tts(text: str, max_chars: int = 200) -> str:
    """清理文字讓 TTS 更口語自然：去除 emoji、Markdown、URL、符號，限制長度"""
    import re

    # 移除 URL
    text = re.sub(r"https?://\S+", "", text)
    text = re.sub(r"www\.\S+", "", text)

    # 移除 emoji（精確範圍，不影響中文字）
    emoji_pattern = re.compile(
        "["
        "\U0001F600-\U0001F64F"
        "\U0001F300-\U0001F5FF"
        "\U0001F680-\U0001F6FF"
        "\U0001F700-\U0001F9FF"
        "\U0001FA00-\U0001FAFF"
        "\u2600-\u26FF"
        "\u2700-\u27BF"
        "\u2300-\u23FF"
        "\u25A0-\u25FF"
        "\u2B00-\u2BFF"
        "\uFE00-\uFE0F"
        "\u200B-\u200D\uFEFF"
        "]+", flags=re.UNICODE
    )
    text = emoji_pattern.sub("", text)

    # 移除 Markdown 格式
    text = re.sub(r"`{1,3}.*?`{1,3}", "", text, flags=re.DOTALL)  # 程式碼塊優先移除
    text = re.sub(r"\*{1,3}(.*?)\*{1,3}", r"\1", text)
    text = re.sub(r"_{1,2}(.*?)_{1,2}", r"\1", text)
    text = re.sub(r"~~(.*?)~~", r"\1", text)
    text = re.sub(r"#{1,6}\s*", "", text)
    text = re.sub(r">\s*", "", text)
    text = re.sub(r"\[([^\]]+)\]\([^\)]+\)", r"\1", text)
    text = re.sub(r"!\[.*?\]\(.*?\)", "", text)

    # 移除表格、程式碼行、路徑
    text = re.sub(r"\|.*?\|", "", text)                          # 表格
    text = re.sub(r"[A-Za-z]:\\[^\s，。！？]*", "", text)        # Windows 路徑
    text = re.sub(r"/[a-zA-Z0-9_/.-]{5,}", "", text)            # Unix 路徑

    # 移除特殊符號
    text = re.sub(r"[|\\/<>{}=+\[\]~^@#$%&*_`]", "", text)

    # 移除列表符號
    text = re.sub(r"^\s*[-•·✅❌]\s*", "", text, flags=re.MULTILINE)
    text = re.sub(r"^\s*\d+[.)]\s*", "", text, flags=re.MULTILINE)

    # 標點正規化
    text = re.sub(r"[!！]{2,}", "！", text)
    text = re.sub(r"[?？]{2,}", "？", text)
    text = re.sub(r"[,，]{2,}", "，", text)
    text = re.sub(r"[.。]{2,}", "。", text)
    text = re.sub(r"\.{2,}", "，", text)

    # 合併空行，去除空白行
    lines = [l.strip() for l in text.splitlines() if l.strip()]
    text = "，".join(lines)

    # 去除多餘空白
    text = re.sub(r"\s{2,}", " ", text).strip()

    # 截斷到 max_chars，在句子邊界截斷
    if len(text) > max_chars:
        cut = text[:max_chars]
        for sep in ["。", "！", "？", "，", " "]:
            idx = cut.rfind(sep)
            if idx > max_chars // 2:
                cut = cut[:idx + 1]
                break
        text = cut

    return text


_XTTS_PORT = 5678
_xtts_server_proc = None

_RVC_PORT = 5679
_rvc_server_proc = None

def _ensure_rvc_server():
    """確保 RVC 伺服器在跑，若沒在跑就啟動它"""
    global _rvc_server_proc
    import urllib.request
    try:
        urllib.request.urlopen(f"http://127.0.0.1:{_RVC_PORT}/health", timeout=1)
        return True
    except Exception:
        pass
    rvc_script = str(Path(__file__).parent / "rvc_server.py")
    if not Path(rvc_script).exists():
        return False
    try:
        _rvc_server_proc = subprocess.Popen(
            [r"C:\Python311\python.exe", rvc_script],
            stdout=subprocess.DEVNULL, stderr=subprocess.DEVNULL
        )
        import time
        for _ in range(30):
            time.sleep(2)
            try:
                urllib.request.urlopen(f"http://127.0.0.1:{_RVC_PORT}/health", timeout=1)
                return True
            except Exception:
                pass
    except Exception as e:
        logging.error(f"RVC 伺服器啟動失敗：{e}")
    return False

def _rvc_convert_wav(wav_bytes: bytes) -> bytes | None:
    """將 WAV bytes 傳給 RVC 伺服器轉換成周杰倫聲音，失敗回傳 None"""
    import urllib.request
    try:
        req = urllib.request.Request(
            f"http://127.0.0.1:{_RVC_PORT}/convert",
            data=wav_bytes,
            headers={"Content-Type": "audio/wav"},
            method="POST"
        )
        with urllib.request.urlopen(req, timeout=60) as resp:
            return resp.read()
    except Exception as e:
        logging.warning(f"RVC 轉換失敗：{e}")
        return None

def _ensure_xtts_server():
    """確保 XTTS 伺服器在跑，若沒在跑就啟動它"""
    global _xtts_server_proc
    import urllib.request
    try:
        urllib.request.urlopen(f"http://127.0.0.1:{_XTTS_PORT}/health", timeout=1)
        return True  # 已在跑
    except Exception:
        pass
    # 啟動 xtts_server.py
    xtts_script = str(Path(__file__).parent / "xtts_server.py")
    if not Path(xtts_script).exists():
        return False
    try:
        _xtts_server_proc = subprocess.Popen(
            [r"C:\Python311\python.exe", xtts_script],
            stdout=subprocess.DEVNULL, stderr=subprocess.DEVNULL,
            creationflags=subprocess.CREATE_NO_WINDOW
        )
        # 等待啟動（模型載入需要 20-30 秒）
        for _ in range(60):
            import time as _t
            _t.sleep(1)
            try:
                urllib.request.urlopen(f"http://127.0.0.1:{_XTTS_PORT}/health", timeout=1)
                return True
            except Exception:
                pass
    except Exception as e:
        logging.error(f"XTTS 伺服器啟動失敗：{e}")
    return False


def _xtts_generate_wav(text: str) -> bytes | None:
    """向 XTTS 伺服器請求語音，回傳 WAV bytes；失敗回傳 None"""
    import json as _json, urllib.request as _req
    payload = _json.dumps({"text": text}).encode("utf-8")
    try:
        req = _req.Request(
            f"http://127.0.0.1:{_XTTS_PORT}/tts",
            data=payload,
            headers={"Content-Type": "application/json"},
            method="POST"
        )
        with _req.urlopen(req, timeout=60) as resp:
            return resp.read()
    except Exception as e:
        logging.warning(f"XTTS 請求失敗：{e}")
        return None


def generate_voice_ogg(text: str, voice: str = "zh-CN-YunxiNeural") -> bytes:
    """生成語音並回傳 OGG OPUS bytes（Telegram voice message 格式）
    優先使用 XTTS v2（更自然人聲），失敗自動 fallback 到 edge_tts
    """
    text = clean_for_tts(text)
    import tempfile, subprocess as sp
    import imageio_ffmpeg

    ffmpeg_exe = imageio_ffmpeg.get_ffmpeg_exe()
    tmp_ogg = tempfile.NamedTemporaryFile(suffix=".ogg", delete=False)
    tmp_ogg.close()

    # ── 嘗試 XTTS v2 ──────────────────────────────────────
    wav_bytes = None
    if _ensure_xtts_server():
        wav_bytes = _xtts_generate_wav(text)

    if wav_bytes:
        # XTTS WAV → OGG OPUS（加低音 EQ）
        tmp_wav = tempfile.NamedTemporaryFile(suffix=".wav", delete=False)
        tmp_wav.write(wav_bytes)
        tmp_wav.close()
        sp.run([
            ffmpeg_exe, "-y", "-i", tmp_wav.name,
            "-af", "equalizer=f=80:width_type=o:width=2:g=6,equalizer=f=150:width_type=o:width=2:g=4",
            "-c:a", "libopus", "-b:a", "96k",
            tmp_ogg.name
        ], capture_output=True)
        Path(tmp_wav.name).unlink(missing_ok=True)
    else:
        # ── Fallback：edge_tts ──────────────────────────────
        import edge_tts, asyncio
        tmp_mp3 = tempfile.NamedTemporaryFile(suffix=".mp3", delete=False)
        tmp_mp3.close()

        async def _gen():
            comm = edge_tts.Communicate(text, voice, rate="-10%", pitch="-15Hz")
            await comm.save(tmp_mp3.name)
        asyncio.run(_gen())

        # MP3 → OGG OPUS（加低音 EQ）
        sp.run([
            ffmpeg_exe, "-y", "-i", tmp_mp3.name,
            "-af", "equalizer=f=80:width_type=o:width=2:g=6,equalizer=f=150:width_type=o:width=2:g=4",
            "-c:a", "libopus", "-b:a", "96k",
            tmp_ogg.name
        ], capture_output=True)
        Path(tmp_mp3.name).unlink(missing_ok=True)

    data = Path(tmp_ogg.name).read_bytes()
    Path(tmp_ogg.name).unlink(missing_ok=True)
    return data


def execute_tts_advanced(action, text="", voice="zh-CN-YunxiNeural"):
    try:
        import edge_tts, asyncio
        if action == "speak":
            out = str(Path("C:/Users/blue_/Desktop/測試檔案") / f"tts_{datetime.datetime.now().strftime('%H%M%S')}.mp3")
            _clean = clean_for_tts(text)
            async def _gen():
                comm = edge_tts.Communicate(_clean, voice, rate="-5%", pitch="-5Hz")
                await comm.save(out)
            asyncio.run(_gen())
            subprocess.Popen(["powershell.exe","-Command",f"Start-Process '{out}'"])
            return f"✅ Edge TTS 語音已播放：{voice}"
        elif action == "list_voices":
            async def _list(): return await edge_tts.list_voices()
            voices = asyncio.run(_list())
            return "\n".join(f"{v['ShortName']}  {v['FriendlyName']}" for v in voices if v["Locale"].startswith("zh"))
    except Exception as e:
        return f"❌ Edge TTS 失敗：{e}"


def execute_todo(action, task="", todo_id=0):
    try:
        db = str(Path("C:/Users/blue_/claude-telegram-bot/todo.db"))
        conn = sqlite3.connect(db)
        conn.execute("CREATE TABLE IF NOT EXISTS todos (id INTEGER PRIMARY KEY AUTOINCREMENT, task TEXT, done INTEGER DEFAULT 0, created TEXT)")
        conn.commit()
        if action == "add":
            conn.execute("INSERT INTO todos (task,created) VALUES (?,?)", (task, datetime.datetime.now().strftime("%Y-%m-%d %H:%M")))
            conn.commit(); conn.close(); return f"✅ 已新增：{task}"
        elif action == "list":
            rows = conn.execute("SELECT id,task,done,created FROM todos ORDER BY done,id").fetchall()
            conn.close()
            return "\n".join(f"{'✅' if r[2] else '⬜'} [{r[0]}] {r[1]}" for r in rows) or "（清單為空）"
        elif action == "done":
            conn.execute("UPDATE todos SET done=1 WHERE id=?", (todo_id,)); conn.commit(); conn.close(); return f"✅ 任務 #{todo_id} 已完成"
        elif action == "delete":
            conn.execute("DELETE FROM todos WHERE id=?", (todo_id,)); conn.commit(); conn.close(); return f"✅ 任務 #{todo_id} 已刪除"
        elif action == "clear":
            conn.execute("DELETE FROM todos WHERE done=1"); conn.commit(); conn.close(); return "✅ 已清除所有已完成任務"
        else:
            conn.close()
            return f"⚠️ 不支援的操作：{action}"
    except Exception as e:
        return f"❌ 任務清單失敗：{e}"


def execute_sysres_chart(duration=10):
    try:
        import psutil, matplotlib
        matplotlib.use("Agg")
        import matplotlib.pyplot as plt, time as t
        cpu_vals, mem_vals = [], []
        for _ in range(duration):
            cpu_vals.append(psutil.cpu_percent(interval=1))
            mem_vals.append(psutil.virtual_memory().percent)
        out = str(Path("C:/Users/blue_/Desktop/測試檔案") / f"sysres_{datetime.datetime.now().strftime('%H%M%S')}.png")
        fig, ax = plt.subplots()
        ax.plot(range(1,duration+1), cpu_vals, label="CPU %", color="blue")
        ax.plot(range(1,duration+1), mem_vals, label="RAM %", color="orange")
        ax.set_ylim(0,100); ax.legend(); ax.set_title("系統資源使用率")
        plt.tight_layout(); plt.savefig(out); plt.close()
        return out
    except Exception as e:
        return f"❌ 失敗：{e}"


def execute_password_mgr(action, site, master, username="", password=""):
    try:
        from cryptography.hazmat.primitives.kdf.pbkdf2 import PBKDF2HMAC
        from cryptography.hazmat.primitives import hashes
        from cryptography.fernet import Fernet
        import base64
        kdf = PBKDF2HMAC(algorithm=hashes.SHA256(), length=32, salt=b"pwd_manager_v1", iterations=100000)
        key = base64.urlsafe_b64encode(kdf.derive(master.encode()))
        f = Fernet(key)
        db = str(Path("C:/Users/blue_/claude-telegram-bot/passwords.db"))
        conn = sqlite3.connect(db)
        conn.execute("CREATE TABLE IF NOT EXISTS passwords (id INTEGER PRIMARY KEY AUTOINCREMENT, site TEXT, username TEXT, password TEXT)")
        conn.commit()
        if action == "save":
            enc = f.encrypt(password.encode()).decode()
            conn.execute("INSERT INTO passwords (site,username,password) VALUES (?,?,?)", (site, username, enc))
            conn.commit(); conn.close(); return f"✅ 密碼已儲存：{site}"
        elif action == "get":
            rows = conn.execute("SELECT site,username,password FROM passwords WHERE site LIKE ?", (f"%{site}%",)).fetchall()
            conn.close()
            if not rows: return f"（找不到 {site} 的密碼）"
            results = []
            for s, u, enc_p in rows:
                try:
                    decrypted = f.decrypt(enc_p.encode()).decode()
                    masked = decrypted[:2] + "*" * max(0, len(decrypted) - 4) + decrypted[-2:] if len(decrypted) > 4 else "****"
                    results.append(f"🔑 {s}\n帳號：{u}\n密碼：{masked}")
                except Exception: results.append(f"❌ {s} 解密失敗（主密碼錯誤？）")
            return "\n".join(results)
        else:
            conn.close()
            return f"⚠️ 不支援的操作：{action}（支援 save/get）"
    except Exception as e:
        return f"❌ 密碼管理失敗：{e}"


def execute_clipboard_image(action, path=""):
    try:
        import win32clipboard
        from PIL import Image
        import io as _io
        if action == "get":
            win32clipboard.OpenClipboard()
            try: data = win32clipboard.GetClipboardData(win32clipboard.CF_DIB)
            finally: win32clipboard.CloseClipboard()
            out = path or str(Path("C:/Users/blue_/Desktop/測試檔案") / f"clipboard_{datetime.datetime.now().strftime('%H%M%S')}.png")
            Image.open(_io.BytesIO(data)).save(out)
            return f"✅ 剪貼簿圖片已存：{out}"
        elif action == "set":
            img = Image.open(path).convert("RGB")
            buf = _io.BytesIO(); img.save(buf,"BMP"); data = buf.getvalue()[14:]
            win32clipboard.OpenClipboard()
            try: win32clipboard.EmptyClipboard(); win32clipboard.SetClipboardData(win32clipboard.CF_DIB, data)
            finally: win32clipboard.CloseClipboard()
            return f"✅ 圖片已複製到剪貼簿"
    except Exception as e:
        return f"❌ 剪貼簿圖片失敗：{e}"


def execute_volume(action, level=None):
    try:
        from ctypes import cast, POINTER
        from comtypes import CLSCTX_ALL
        from pycaw.pycaw import AudioUtilities, IAudioEndpointVolume
        devices = AudioUtilities.GetSpeakers()
        interface = devices.Activate(IAudioEndpointVolume._iid_, CLSCTX_ALL, None)
        volume = cast(interface, POINTER(IAudioEndpointVolume))
        if action == "get":
            vol = int(volume.GetMasterVolumeLevelScalar() * 100)
            muted = bool(volume.GetMute())
            return f"🔊 音量：{vol}%{'（靜音中）' if muted else ''}"
        elif action == "set":
            volume.SetMasterVolumeLevelScalar(max(0, min(100, int(level))) / 100.0, None)
            return f"✅ 音量設定為 {level}%"
        elif action == "mute":
            volume.SetMute(1, None)
            return "✅ 已靜音"
        elif action == "unmute":
            volume.SetMute(0, None)
            return "✅ 已取消靜音"
    except Exception as e:
        return f"❌ 音量控制失敗：{e}"

def execute_display(action, level=None):
    try:
        import screen_brightness_control as sbc
        if action == "brightness_get":
            b = sbc.get_brightness()
            return f"💡 亮度：{b}%"
        elif action == "brightness_set":
            sbc.set_brightness(max(0, min(100, int(level))))
            return f"✅ 亮度設定為 {level}%"
        elif action == "resolution":
            import subprocess
            result = subprocess.run(["powershell", "-Command",
                "Get-CimInstance Win32_VideoController | Select-Object CurrentHorizontalResolution,CurrentVerticalResolution,VideoModeDescription | Format-List"],
                capture_output=True, text=True)
            return f"🖥️ 解析度資訊：\n{result.stdout.strip()}"
    except Exception as e:
        return f"❌ 螢幕控制失敗：{e}"

def execute_media(action, device_name=""):
    try:
        import keyboard
        key_map = {
            "play_pause": "play/pause media",
            "next": "next track",
            "prev": "previous track",
            "stop": "stop media",
            "volume_up": "volume up",
            "volume_down": "volume down",
            "mute": "volume mute"
        }
        if action in key_map:
            keyboard.send(key_map[action])
            return f"✅ 媒體控制：{action}"
        elif action == "list_devices":
            from pycaw.pycaw import AudioUtilities
            devices = AudioUtilities.GetAllDevices()
            names = [d.FriendlyName for d in devices if d.FriendlyName]
            return "🔊 音訊裝置：\n" + "\n".join(f"- {n}" for n in names)
        elif action == "switch_device":
            import subprocess
            result = subprocess.run(["powershell", "-Command",
                f"Get-AudioDevice -List | Where-Object {{$_.Name -like '*{device_name}*'}} | Set-AudioDevice"],
                capture_output=True, text=True)
            if result.returncode == 0:
                return f"✅ 已切換至：{device_name}"
            return f"⚠️ 切換失敗（需安裝 AudioDeviceCmdlets）：{result.stderr}"
    except Exception as e:
        return f"❌ 媒體控制失敗：{e}"

def execute_software(action, name="", keyword=""):
    try:
        import subprocess
        if action == "list":
            q = f"| Where-Object {{$_.DisplayName -like '*{keyword}*'}}" if keyword else ""
            result = subprocess.run(["powershell", "-Command",
                f"Get-ItemProperty HKLM:\\Software\\Wow6432Node\\Microsoft\\Windows\\CurrentVersion\\Uninstall\\* {q} | Select-Object DisplayName,DisplayVersion | Sort-Object DisplayName | Format-Table -AutoSize"],
                capture_output=True, text=True, encoding="utf-8", errors="replace")
            return f"📦 已安裝軟體：\n{result.stdout.strip()[:2000]}"
        elif action == "install":
            result = subprocess.run(["winget", "install", "--id", name, "-e", "--silent"],
                capture_output=True, text=True)
            if result.returncode == 0:
                return f"✅ 安裝成功：{name}"
            return f"⚠️ 安裝結果：{result.stdout[-500:]}"
        elif action == "uninstall":
            result = subprocess.run(["winget", "uninstall", "--name", name, "--silent"],
                capture_output=True, text=True)
            if result.returncode == 0:
                return f"✅ 已卸載：{name}"
            return f"⚠️ 卸載結果：{result.stdout[-500:]}"
    except Exception as e:
        return f"❌ 軟體管理失敗：{e}"

def execute_startup(action, name="", command=""):
    try:
        import subprocess, winreg
        key_path = r"Software\Microsoft\Windows\CurrentVersion\Run"
        if action == "list":
            result = subprocess.run(["powershell", "-Command",
                "Get-CimInstance Win32_StartupCommand | Select-Object Name,Command,Location | Format-Table -AutoSize"],
                capture_output=True, text=True, encoding="utf-8", errors="replace")
            return f"🚀 開機自啟：\n{result.stdout.strip()}"
        elif action == "add":
            with winreg.OpenKey(winreg.HKEY_CURRENT_USER, key_path, 0, winreg.KEY_SET_VALUE) as k:
                winreg.SetValueEx(k, name, 0, winreg.REG_SZ, command)
            return f"✅ 已新增開機自啟：{name}"
        elif action == "remove":
            try:
                with winreg.OpenKey(winreg.HKEY_CURRENT_USER, key_path, 0, winreg.KEY_SET_VALUE) as k:
                    winreg.DeleteValue(k, name)
                return f"✅ 已移除開機自啟：{name}"
            except FileNotFoundError:
                return f"⚠️ 找不到項目：{name}"
    except Exception as e:
        return f"❌ 開機自啟管理失敗：{e}"

def execute_env_var(action, name="", value="", permanent=False):
    try:
        import os, subprocess
        if action == "get":
            v = os.environ.get(name, "")
            if not v:
                result = subprocess.run(["powershell", "-Command",
                    f"[System.Environment]::GetEnvironmentVariable('{name}','Machine')"],
                    capture_output=True, text=True)
                v = result.stdout.strip()
            return f"🌍 {name} = {v}" if v else f"⚠️ 環境變數 {name} 不存在"
        elif action == "set":
            os.environ[name] = value
            if permanent:
                subprocess.run(["powershell", "-Command",
                    f"[System.Environment]::SetEnvironmentVariable('{name}','{value}','User')"],
                    capture_output=True)
            return f"✅ 環境變數設定：{name}={value}{'（永久）' if permanent else '（本次）'}"
    except Exception as e:
        return f"❌ 環境變數失敗：{e}"

def execute_user_account(action, username="", password=""):
    try:
        import subprocess
        if action == "list":
            result = subprocess.run(["powershell", "-Command",
                "Get-LocalUser | Select-Object Name,Enabled,LastLogon | Format-Table -AutoSize"],
                capture_output=True, text=True, encoding="utf-8", errors="replace")
            return f"👤 使用者帳戶：\n{result.stdout.strip()}"
        elif action == "create":
            ps = f"$pw = ConvertTo-SecureString '{password}' -AsPlainText -Force; New-LocalUser '{username}' -Password $pw -FullName '{username}'"
            result = subprocess.run(["powershell", "-Command", ps], capture_output=True, text=True)
            if result.returncode == 0:
                return f"✅ 已建立帳戶：{username}"
            return f"❌ 建立失敗：{result.stderr.strip()}"
        elif action == "delete":
            result = subprocess.run(["powershell", "-Command", f"Remove-LocalUser '{username}'"],
                capture_output=True, text=True)
            if result.returncode == 0:
                return f"✅ 已刪除帳戶：{username}"
            return f"❌ 刪除失敗：{result.stderr.strip()}"
    except Exception as e:
        return f"❌ 使用者帳戶管理失敗：{e}"

def execute_win_update(action):
    try:
        import subprocess
        if action in ("list", "check"):
            result = subprocess.run(["powershell", "-Command",
                "Get-WindowsUpdate -AcceptAll -Verbose 2>&1 | Select-Object -First 20"],
                capture_output=True, text=True, timeout=60)
            out = result.stdout.strip() or "無可用更新或需要 PSWindowsUpdate 模組"
            return f"🔄 Windows Update：\n{out[:1500]}"
        elif action == "install":
            result = subprocess.run(["powershell", "-Command",
                "Install-WindowsUpdate -AcceptAll -AutoReboot:$false -Verbose 2>&1"],
                capture_output=True, text=True, timeout=300)
            return f"✅ 更新執行完成：\n{result.stdout.strip()[:1500]}"
    except subprocess.TimeoutExpired:
        return "⏳ 更新查詢超時，請稍後再試"
    except Exception as e:
        return f"❌ Windows Update 失敗：{e}"

def execute_device_manager(action, name="", keyword=""):
    try:
        import subprocess
        if action == "list":
            q = f"| Where-Object {{$_.Name -like '*{keyword}*'}}" if keyword else ""
            result = subprocess.run(["powershell", "-Command",
                f"Get-PnpDevice {q} | Select-Object Status,Class,FriendlyName | Format-Table -AutoSize"],
                capture_output=True, text=True, encoding="utf-8", errors="replace")
            return f"🖥️ 裝置清單：\n{result.stdout.strip()[:2000]}"
        elif action == "enable":
            result = subprocess.run(["powershell", "-Command",
                f"Get-PnpDevice | Where-Object {{$_.FriendlyName -like '*{name}*'}} | Enable-PnpDevice -Confirm:$false"],
                capture_output=True, text=True)
            return f"✅ 已啟用裝置：{name}" if result.returncode == 0 else f"❌ 啟用失敗：{result.stderr}"
        elif action == "disable":
            result = subprocess.run(["powershell", "-Command",
                f"Get-PnpDevice | Where-Object {{$_.FriendlyName -like '*{name}*'}} | Disable-PnpDevice -Confirm:$false"],
                capture_output=True, text=True)
            return f"✅ 已停用裝置：{name}" if result.returncode == 0 else f"❌ 停用失敗：{result.stderr}"
    except Exception as e:
        return f"❌ 裝置管理員失敗：{e}"

def execute_network_config(action, name="", ip="", dns1="", dns2="", domain="", duration=10):
    try:
        import subprocess, time
        if action == "adapter_list":
            result = subprocess.run(["powershell", "-Command",
                "Get-NetAdapter | Select-Object Name,Status,LinkSpeed | Format-Table -AutoSize"],
                capture_output=True, text=True, encoding="utf-8", errors="replace")
            return f"🌐 網路介面卡：\n{result.stdout.strip()}"
        elif action == "adapter_enable":
            subprocess.run(["powershell", "-Command", f"Enable-NetAdapter -Name '{name}' -Confirm:$false"], capture_output=True)
            return f"✅ 已啟用介面卡：{name}"
        elif action == "adapter_disable":
            subprocess.run(["powershell", "-Command", f"Disable-NetAdapter -Name '{name}' -Confirm:$false"], capture_output=True)
            return f"✅ 已停用介面卡：{name}"
        elif action == "dns_get":
            result = subprocess.run(["powershell", "-Command",
                f"Get-DnsClientServerAddress -InterfaceAlias '{name}' | Format-List"],
                capture_output=True, text=True, encoding="utf-8", errors="replace")
            return f"🌐 DNS 設定：\n{result.stdout.strip()}"
        elif action == "dns_set":
            servers = f"'{dns1}','{dns2}'" if dns2 else f"'{dns1}'"
            subprocess.run(["powershell", "-Command",
                f"Set-DnsClientServerAddress -InterfaceAlias '{name}' -ServerAddresses ({servers})"],
                capture_output=True)
            return f"✅ DNS 設定完成：{dns1}" + (f", {dns2}" if dns2 else "")
        elif action == "ip_set":
            subprocess.run(["powershell", "-Command",
                f"New-NetIPAddress -InterfaceAlias '{name}' -IPAddress '{ip}' -PrefixLength 24 -DefaultGateway '{dns1}' -Confirm:$false 2>&1"],
                capture_output=True)
            return f"✅ 靜態 IP 設定：{ip}"
        elif action == "hosts_list":
            with open(r"C:\Windows\System32\drivers\etc\hosts", "r", encoding="utf-8", errors="replace") as f:
                content = f.read()
            return f"📄 Hosts 檔案：\n{content[:1500]}"
        elif action == "hosts_add":
            entry = f"\n{ip}\t{domain}"
            with open(r"C:\Windows\System32\drivers\etc\hosts", "a", encoding="utf-8") as f:
                f.write(entry)
            return f"✅ 已新增 hosts：{ip} → {domain}"
        elif action == "hosts_remove":
            with open(r"C:\Windows\System32\drivers\etc\hosts", "r", encoding="utf-8", errors="replace") as f:
                lines = f.readlines()
            new_lines = [l for l in lines if domain not in l]
            with open(r"C:\Windows\System32\drivers\etc\hosts", "w", encoding="utf-8") as f:
                f.writelines(new_lines)
            return f"✅ 已移除 hosts：{domain}"
        elif action == "traffic":
            import psutil
            t1 = psutil.net_io_counters()
            time.sleep(min(int(duration), 10))
            t2 = psutil.net_io_counters()
            sent = (t2.bytes_sent - t1.bytes_sent) / 1024
            recv = (t2.bytes_recv - t1.bytes_recv) / 1024
            return f"📊 網路流量（{duration}s）：↑ {sent:.1f} KB  ↓ {recv:.1f} KB"
    except Exception as e:
        return f"❌ 網路設定失敗：{e}"

def execute_automation(action, condition_type="", condition_value="", command="",
                       duration=60.0, layout="side_by_side", x=0, y=0, w=0, h=0,
                       keyword="", output=""):
    try:
        if action == "if_then":
            import subprocess, psutil, time, os
            deadline = time.time() + float(duration)
            while time.time() < deadline:
                triggered = False
                if condition_type == "cpu_above":
                    triggered = psutil.cpu_percent(1) > float(condition_value)
                elif condition_type == "mem_above":
                    triggered = psutil.virtual_memory().percent > float(condition_value)
                elif condition_type == "file_exists":
                    triggered = Path(condition_value).exists()
                elif condition_type == "process_running":
                    triggered = any(condition_value.lower() in p.name().lower() for p in psutil.process_iter())
                elif condition_type == "time_is":
                    triggered = datetime.datetime.now().strftime("%H:%M") == condition_value
                if triggered:
                    subprocess.Popen(command, shell=True)
                    return f"✅ 條件達成（{condition_type}={condition_value}），已執行：{command}"
                time.sleep(2)
            return f"⏳ 監控 {duration}s 內條件未達成"
        elif action == "window_arrange":
            import win32gui, win32con
            layouts = {"side_by_side": [], "quad": [], "stack": [], "maximize_all": []}
            import ctypes
            user32 = ctypes.windll.user32
            sw = user32.GetSystemMetrics(0)
            sh = user32.GetSystemMetrics(1)
            hwnds = []
            win32gui.EnumWindows(lambda h, l: l.append(h) if win32gui.IsWindowVisible(h) and win32gui.GetWindowText(h) else None, hwnds)
            if layout == "side_by_side" and len(hwnds) >= 2:
                win32gui.MoveWindow(hwnds[0], 0, 0, sw//2, sh, True)
                win32gui.MoveWindow(hwnds[1], sw//2, 0, sw//2, sh, True)
                return f"✅ 左右排列完成（{sw}x{sh}）"
            elif layout == "quad" and len(hwnds) >= 4:
                positions = [(0,0,sw//2,sh//2),(sw//2,0,sw//2,sh//2),(0,sh//2,sw//2,sh//2),(sw//2,sh//2,sw//2,sh//2)]
                for i, (x_,y_,w_,h_) in enumerate(positions[:4]):
                    win32gui.MoveWindow(hwnds[i], x_, y_, w_, h_, True)
                return "✅ 四象限排列完成"
            elif layout == "stack":
                h_each = sh // max(len(hwnds), 1)
                for i, hwnd in enumerate(hwnds[:8]):
                    win32gui.MoveWindow(hwnd, 0, i*h_each, sw, h_each, True)
                return f"✅ 堆疊排列完成（{len(hwnds[:8])} 個視窗）"
            elif layout == "maximize_all":
                for hwnd in hwnds:
                    win32gui.ShowWindow(hwnd, win32con.SW_MAXIMIZE)
                return f"✅ 全部最大化（{len(hwnds)} 個）"
            return f"✅ 排列完成"
        elif action == "region_ocr":
            import pyautogui, easyocr
            region = (int(x), int(y), int(w), int(h)) if w and h else None
            screenshot = pyautogui.screenshot(region=region)
            out = output or str(Path("C:/Users/blue_/Desktop/測試檔案") / "region_ocr.png")
            screenshot.save(out)
            reader = easyocr.Reader(["ch_tra","en"], gpu=False)
            results = reader.readtext(out, detail=0)
            return f"🔍 區域 OCR 結果：\n" + "\n".join(results)
        elif action == "window_screenshot":
            import win32gui, win32ui, win32con, ctypes
            from PIL import Image
            import array
            hwnds = []
            win32gui.EnumWindows(
                lambda h, l: l.append(h) if win32gui.IsWindowVisible(h) and keyword.lower() in win32gui.GetWindowText(h).lower() else None,
                hwnds)
            if not hwnds:
                return f"⚠️ 找不到視窗：{keyword}"
            hwnd = hwnds[0]
            win32gui.SetForegroundWindow(hwnd)
            rect = win32gui.GetWindowRect(hwnd)
            x_, y_, x2, y2 = rect
            w_, h_ = x2-x_, y2-y_
            import pyautogui
            screenshot = pyautogui.screenshot(region=(x_, y_, w_, h_))
            out = output or str(Path("C:/Users/blue_/Desktop/測試檔案") / f"window_{keyword}_{datetime.datetime.now().strftime('%H%M%S')}.png")
            screenshot.save(out)
            return f"✅ 視窗截圖：{out}"
    except Exception as e:
        return f"❌ 自動化失敗：{e}"


_alert_monitors = {}
_interval_schedules = {}
_clipboard_hist = []
_voice_cmd_running = False

def execute_vision_loop(goal, max_steps=20, interval=3.0, timeout=120.0):
    try:
        import pyautogui, anthropic, base64, io, time
        steps = 0
        start = time.time()
        log = []
        _client = anthropic.Anthropic(api_key=os.getenv("ANTHROPIC_API_KEY"))
        while steps < int(max_steps) and (time.time() - start) < float(timeout):
            screenshot = pyautogui.screenshot()
            buf = io.BytesIO()
            screenshot.save(buf, format="PNG")
            img_b64 = base64.standard_b64encode(buf.getvalue()).decode()
            resp = _client.messages.create(
                model="claude-sonnet-4-6",
                max_tokens=512,
                messages=[{
                    "role": "user",
                    "content": [
                        {"type": "image", "source": {"type": "base64", "media_type": "image/png", "data": img_b64}},
                        {"type": "text", "text": f"目標：{goal}\n已執行步驟：{log}\n\n請分析畫面，回答：\n1. 目標是否已達成？（是/否）\n2. 如果否，下一步應該怎麼做？請用 JSON 格式回答：{{\"done\": true/false, \"action\": \"動作說明\", \"type\": \"click/type/key/wait\", \"x\": 0, \"y\": 0, \"text\": \"\"}}"}
                    ]
                }]
            )
            import json, re
            text = resp.content[0].text
            m = re.search(r'\{.*\}', text, re.DOTALL)
            if not m:
                log.append(f"步驟{steps+1}: AI回應無法解析")
                steps += 1; time.sleep(float(interval)); continue
            action = json.loads(m.group())
            if action.get("done"):
                return f"✅ 目標達成！共執行 {steps} 步\n" + "\n".join(log)
            act_type = action.get("type","")
            act_desc = action.get("action","")
            if act_type == "click" and action.get("x"):
                pyautogui.click(action["x"], action["y"])
            elif act_type == "type" and action.get("text"):
                pyautogui.typewrite(action["text"], interval=0.05)
            elif act_type == "key" and action.get("text"):
                pyautogui.press(action["text"])
            elif act_type == "wait":
                time.sleep(2)
            log.append(f"步驟{steps+1}: {act_desc}")
            steps += 1
            time.sleep(float(interval))
        return f"⏳ 達到上限（{steps} 步 / {int(time.time()-start)}s）\n執行記錄：\n" + "\n".join(log)
    except Exception as e:
        return f"❌ 視覺自動化循環失敗：{e}"


def execute_alert_monitor(action, name="", condition="", threshold="", target="", interval=30, chat_id=None, _bot_send=None):
    global _alert_monitors
    try:
        import threading, time, psutil
        if action == "list":
            if not _alert_monitors:
                return "⚠️ 無執行中的監控"
            return "📊 監控清單：\n" + "\n".join(f"- {k}: {v['condition']} {v['threshold']}" for k,v in _alert_monitors.items())
        elif action == "stop":
            if name in _alert_monitors:
                _alert_monitors[name]["running"] = False
                del _alert_monitors[name]
                return f"✅ 已停止監控：{name}"
            return f"⚠️ 找不到監控：{name}"
        elif action == "start":
            if name in _alert_monitors:
                return f"⚠️ 已有同名監控：{name}"
            cfg = {"condition": condition, "threshold": threshold, "target": target, "running": True}
            _alert_monitors[name] = cfg
            send_chat_id = chat_id or OWNER_ID
            def _monitor():
                import easyocr
                reader = None
                while _alert_monitors.get(name, {}).get("running"):
                    try:
                        triggered = False
                        msg = ""
                        val = threshold
                        if condition == "cpu_above":
                            v = psutil.cpu_percent(1)
                            if v > float(val): triggered = True; msg = f"⚠️ CPU 使用率 {v:.1f}% 超過 {val}%"
                        elif condition == "mem_above":
                            v = psutil.virtual_memory().percent
                            if v > float(val): triggered = True; msg = f"⚠️ 記憶體使用率 {v:.1f}% 超過 {val}%"
                        elif condition == "disk_above":
                            v = psutil.disk_usage("/").percent
                            if v > float(val): triggered = True; msg = f"⚠️ 磁碟使用率 {v:.1f}% 超過 {val}%"
                        elif condition == "process_missing":
                            pnames = [p.name().lower() for p in psutil.process_iter(["name"])]
                            if not any(target.lower() in n for n in pnames): triggered = True; msg = f"⚠️ 程序 {target} 已停止執行"
                        elif condition == "process_running":
                            pnames = [p.name().lower() for p in psutil.process_iter(["name"])]
                            if any(target.lower() in n for n in pnames): triggered = True; msg = f"ℹ️ 程序 {target} 正在執行"
                        elif condition == "screen_text_found":
                            import pyautogui
                            if reader is None: reader = easyocr.Reader(["ch_tra","en"], gpu=False)
                            screenshot = pyautogui.screenshot()
                            import tempfile; tmp = tempfile.mktemp(suffix=".png"); screenshot.save(tmp)
                            results = reader.readtext(tmp, detail=0)
                            full_text = " ".join(results)
                            Path(tmp).unlink(missing_ok=True)
                            if target.lower() in full_text.lower(): triggered = True; msg = f"ℹ️ 螢幕偵測到文字：{target}"
                        if triggered and _bot_send:
                            import asyncio
                            asyncio.run_coroutine_threadsafe(_bot_send(send_chat_id, f"🔔 【監控告警】{name}\n{msg}"), _bot_send.__self__.loop if hasattr(_bot_send, '__self__') else asyncio.get_event_loop())
                    except Exception:
                        pass
                    time.sleep(int(interval))
            threading.Thread(target=_monitor, daemon=True).start()
            return f"✅ 監控已啟動：{name}（{condition} {threshold}，每 {interval}s 檢查）"
    except Exception as e:
        return f"❌ 監控告警失敗：{e}"


def execute_interval_schedule(action, name="", command="", every_minutes=60.0, repeat=0, duration_hours=0.0):
    global _interval_schedules
    try:
        import threading, time, subprocess
        if action == "list":
            if not _interval_schedules:
                return "⚠️ 無執行中排程"
            return "⏱️ 間隔排程：\n" + "\n".join(f"- {k}: 每{v['mins']}分鐘，已執行{v['count']}次" for k,v in _interval_schedules.items())
        elif action == "stop":
            if name in _interval_schedules:
                _interval_schedules[name]["running"] = False
                del _interval_schedules[name]
                return f"✅ 已停止排程：{name}"
            return f"⚠️ 找不到排程：{name}"
        elif action == "start":
            if name in _interval_schedules:
                return f"⚠️ 已有同名排程：{name}"
            cfg = {"command": command, "mins": every_minutes, "repeat": repeat, "count": 0, "running": True}
            _interval_schedules[name] = cfg
            def _run():
                import time as t
                end_time = t.time() + float(duration_hours) * 3600 if duration_hours > 0 else float("inf")
                max_count = int(repeat) if repeat > 0 else float("inf")
                while _interval_schedules.get(name, {}).get("running"):
                    if t.time() > end_time or _interval_schedules.get(name, {}).get("count", 0) >= max_count:
                        _interval_schedules.pop(name, None); break
                    subprocess.Popen(command, shell=True)
                    _interval_schedules[name]["count"] = _interval_schedules[name].get("count", 0) + 1
                    t.sleep(float(every_minutes) * 60)
            threading.Thread(target=_run, daemon=True).start()
            desc = f"每 {every_minutes} 分鐘" + (f"，共 {repeat} 次" if repeat else "") + (f"，持續 {duration_hours} 小時" if duration_hours else "")
            return f"✅ 間隔排程已啟動：{name}（{desc}）"
    except Exception as e:
        return f"❌ 間隔排程失敗：{e}"


def execute_wait_for_text(text, timeout=60.0, interval=2.0, region=""):
    try:
        import pyautogui, easyocr, time, tempfile
        reader = easyocr.Reader(["ch_tra","en"], gpu=False)
        start = time.time()
        reg = None
        if region:
            parts = [int(v) for v in region.split(",")]
            if len(parts) == 4: reg = tuple(parts)
        while time.time() - start < float(timeout):
            screenshot = pyautogui.screenshot(region=reg)
            tmp = tempfile.mktemp(suffix=".png")
            screenshot.save(tmp)
            results = reader.readtext(tmp, detail=0)
            Path(tmp).unlink(missing_ok=True)
            full = " ".join(results)
            if text.lower() in full.lower():
                elapsed = time.time() - start
                return f"✅ 偵測到文字「{text}」（等待 {elapsed:.1f}s）"
            time.sleep(float(interval))
        return f"⏳ 超時（{timeout}s），未偵測到文字「{text}」"
    except Exception as e:
        return f"❌ 等待文字失敗：{e}"


_browser_page = None

def execute_browser_advanced(action, selector="", value="", name="", tab_index=0, timeout=30.0, url_pattern=""):
    global _browser_page
    try:
        from playwright.sync_api import sync_playwright
        import json
        if _browser_page is None or _browser_page.is_closed():
            return "⚠️ 瀏覽器未開啟，請先用 browser_control 開啟瀏覽器"
        page = _browser_page
        if action == "wait_element":
            page.wait_for_selector(selector, timeout=float(timeout)*1000)
            return f"✅ 元素已出現：{selector}"
        elif action == "switch_frame":
            frame = page.frame(selector) or page.frame_locator(selector).first
            return f"✅ 已切換 iframe：{selector}"
        elif action == "get_cookies":
            cookies = page.context.cookies()
            lines = [f"{c['name']}={c['value'][:30]}" for c in cookies[:20]]
            return "🍪 Cookies：\n" + "\n".join(lines)
        elif action == "set_cookie":
            page.context.add_cookies([{"name": name, "value": value, "url": page.url}])
            return f"✅ Cookie 已設定：{name}={value}"
        elif action == "list_tabs":
            pages = page.context.pages
            lines = [f"{i}: {p.title()} - {p.url[:50]}" for i, p in enumerate(pages)]
            return "📑 所有分頁：\n" + "\n".join(lines)
        elif action == "switch_tab":
            pages = page.context.pages
            if tab_index < len(pages):
                pages[tab_index].bring_to_front()
                _browser_page = pages[tab_index]
                return f"✅ 已切換到分頁 {tab_index}：{pages[tab_index].title()}"
            return f"⚠️ 分頁 {tab_index} 不存在"
        elif action == "new_tab":
            new_page = page.context.new_page()
            if value: new_page.goto(value)
            _browser_page = new_page
            return f"✅ 已開新分頁：{value or '空白'}"
        elif action == "close_tab":
            page.close()
            pages = page.context.pages
            if pages: _browser_page = pages[-1]
            return "✅ 已關閉目前分頁"
        elif action == "fill_form":
            import json as _json
            fields = _json.loads(value) if value.startswith("{") else {}
            for sel, val in fields.items():
                page.fill(sel, val)
            return f"✅ 表單已填寫"
        elif action == "select_option":
            page.select_option(selector, value)
            return f"✅ 已選擇：{value}"
        elif action == "scroll_to":
            page.locator(selector).scroll_into_view_if_needed()
            return f"✅ 已滾動到：{selector}"
        elif action == "get_html":
            return page.inner_html(selector or "body")[:2000]
        elif action == "wait_url":
            page.wait_for_url(f"**{url_pattern}**", timeout=float(timeout)*1000)
            return f"✅ URL 已包含：{url_pattern}"
    except Exception as e:
        return f"❌ 瀏覽器進階操作失敗：{e}"


_voice_cmd_running = False

def execute_voice_cmd(action, duration=300.0, language="zh-TW", _bot_send=None, _chat_id=None):
    global _voice_cmd_running
    try:
        import threading
        if action == "stop":
            _voice_cmd_running = False
            return "✅ 語音命令模式已停止"
        elif action == "start":
            if _voice_cmd_running:
                return "⚠️ 語音命令模式已在執行中"
            _voice_cmd_running = True
            def _listen_loop():
                global _voice_cmd_running
                import sounddevice as sd, soundfile as sf, speech_recognition as sr, tempfile, time, subprocess
                recognizer = sr.Recognizer()
                sample_rate = 16000
                end_time = time.time() + float(duration)
                while _voice_cmd_running and time.time() < end_time:
                    try:
                        recording = sd.rec(int(4 * sample_rate), samplerate=sample_rate, channels=1, dtype="int16")
                        sd.wait()
                        with tempfile.NamedTemporaryFile(suffix=".wav", delete=False) as f:
                            tmp = f.name
                        sf.write(tmp, recording, sample_rate)
                        with sr.AudioFile(tmp) as source:
                            audio = recognizer.record(source)
                        Path(tmp).unlink(missing_ok=True)
                        text = recognizer.recognize_google(audio, language=language)
                        if not text: continue
                        if "停止" in text or "stop" in text.lower():
                            _voice_cmd_running = False
                            if _bot_send and _chat_id:
                                import asyncio
                                asyncio.run_coroutine_threadsafe(_bot_send(_chat_id, "🎤 語音命令模式已停止"), asyncio.get_event_loop())
                            break
                        if _bot_send and _chat_id:
                            import asyncio
                            asyncio.run_coroutine_threadsafe(_bot_send(_chat_id, f"🎤 語音命令：{text}"), asyncio.get_event_loop())
                        subprocess.Popen(text, shell=True)
                    except Exception:
                        pass
                _voice_cmd_running = False
            threading.Thread(target=_listen_loop, daemon=True).start()
            return f"✅ 語音命令模式已啟動（{duration}s），說「停止」結束"
    except Exception as e:
        return f"❌ 語音命令模式失敗：{e}"


def execute_win_notify_relay(action, duration=3600.0, filter_app="", _bot_send=None, _chat_id=None):
    try:
        import threading, time, subprocess
        if action == "status":
            return "ℹ️ Windows 通知攔截透過輪詢事件記錄實現"
        elif action == "stop":
            return "✅ 通知攔截已標記停止（重啟生效）"
        elif action == "start":
            def _relay():
                import win32evtlog, win32evtlogutil, time as t
                seen = set()
                end = t.time() + float(duration)
                hand = win32evtlog.OpenEventLog(None, "Application")
                while t.time() < end:
                    try:
                        flags = win32evtlog.EVENTLOG_BACKWARDS_READ | win32evtlog.EVENTLOG_SEQUENTIAL_READ
                        batch = win32evtlog.ReadEventLog(hand, flags, 0)
                        for e in (batch or []):
                            eid = (e.RecordNumber, e.TimeGenerated.Format())
                            if eid in seen: continue
                            seen.add(eid)
                            src = e.SourceName
                            if filter_app and filter_app.lower() not in src.lower(): continue
                            try: msg = win32evtlogutil.SafeFormatMessage(e, "Application")[:200]
                            except Exception: msg = "(無法讀取)"
                            if _bot_send and _chat_id and e.EventType in (1, 2):
                                import asyncio
                                level = "❌" if e.EventType == 1 else "⚠️"
                                asyncio.run_coroutine_threadsafe(
                                    _bot_send(_chat_id, f"🔔 Windows 通知\n{level} {src}\n{msg}"),
                                    asyncio.get_event_loop())
                    except Exception:
                        pass
                    t.sleep(10)
                win32evtlog.CloseEventLog(hand)
            threading.Thread(target=_relay, daemon=True).start()
            return f"✅ Windows 通知攔截已啟動（{duration}s）" + (f"，過濾：{filter_app}" if filter_app else "")
    except Exception as e:
        return f"❌ 通知攔截失敗：{e}"


def execute_data_process(action, path="", output="", query="", data="", paths=""):
    try:
        import json, csv, io as _io
        if action == "read_json":
            content = json.loads(Path(path).read_text(encoding="utf-8"))
            if isinstance(content, list):
                return f"📄 JSON（{len(content)} 筆）：\n" + json.dumps(content[:5], ensure_ascii=False, indent=2)
            return f"📄 JSON：\n" + json.dumps(content, ensure_ascii=False, indent=2)[:2000]
        elif action == "write_json":
            obj = json.loads(data)
            Path(output or path).write_text(json.dumps(obj, ensure_ascii=False, indent=2), encoding="utf-8")
            return f"✅ JSON 已儲存：{output or path}"
        elif action == "read_csv":
            rows = []
            with open(path, encoding="utf-8-sig", errors="replace") as f:
                reader = csv.DictReader(f)
                rows = list(reader)
            return f"📊 CSV（{len(rows)} 筆，欄位：{list(rows[0].keys()) if rows else []}）：\n" + json.dumps(rows[:5], ensure_ascii=False, indent=2)
        elif action == "write_csv":
            obj = json.loads(data)
            if not obj: return "⚠️ 資料為空"
            out = output or path
            with open(out, "w", newline="", encoding="utf-8-sig") as f:
                writer = csv.DictWriter(f, fieldnames=obj[0].keys())
                writer.writeheader(); writer.writerows(obj)
            return f"✅ CSV 已儲存：{out}"
        elif action == "filter":
            with open(path, encoding="utf-8-sig", errors="replace") as f:
                rows = list(csv.DictReader(f))
            filtered = []
            for row in rows:
                try:
                    if eval(query, {"__builtins__": {}}, row): filtered.append(row)
                except Exception: pass
            return f"📊 過濾結果（{len(filtered)}/{len(rows)} 筆）：\n" + json.dumps(filtered[:10], ensure_ascii=False, indent=2)
        elif action == "stats":
            with open(path, encoding="utf-8-sig", errors="replace") as f:
                rows = list(csv.DictReader(f))
            if not rows: return "⚠️ 無資料"
            cols = rows[0].keys()
            stats = {}
            for col in cols:
                vals = [row[col] for row in rows if row[col]]
                try:
                    nums = [float(v) for v in vals]
                    stats[col] = {"count": len(nums), "min": min(nums), "max": max(nums), "avg": round(sum(nums)/len(nums),2)}
                except Exception:
                    stats[col] = {"count": len(vals), "unique": len(set(vals))}
            return "📊 統計：\n" + json.dumps(stats, ensure_ascii=False, indent=2)
        elif action == "convert":
            ext_in = Path(path).suffix.lower()
            ext_out = Path(output).suffix.lower() if output else ""
            if ext_in == ".json" and ext_out == ".csv":
                obj = json.loads(Path(path).read_text(encoding="utf-8"))
                if not isinstance(obj, list): obj = [obj]
                with open(output, "w", newline="", encoding="utf-8-sig") as f:
                    w = csv.DictWriter(f, fieldnames=obj[0].keys()); w.writeheader(); w.writerows(obj)
                return f"✅ JSON → CSV：{output}"
            elif ext_in == ".csv" and ext_out == ".json":
                with open(path, encoding="utf-8-sig", errors="replace") as f:
                    rows = list(csv.DictReader(f))
                Path(output).write_text(json.dumps(rows, ensure_ascii=False, indent=2), encoding="utf-8")
                return f"✅ CSV → JSON：{output}"
            return f"⚠️ 不支援轉換：{ext_in} → {ext_out}"
        elif action == "merge":
            all_rows = []
            for p in paths.split(","):
                p = p.strip()
                if p.endswith(".csv"):
                    with open(p, encoding="utf-8-sig", errors="replace") as f:
                        all_rows.extend(list(csv.DictReader(f)))
                elif p.endswith(".json"):
                    obj = json.loads(Path(p).read_text(encoding="utf-8"))
                    if isinstance(obj, list): all_rows.extend(obj)
            out = output or str(Path(paths.split(",")[0].strip()).parent / "merged.csv")
            with open(out, "w", newline="", encoding="utf-8-sig") as f:
                w = csv.DictWriter(f, fieldnames=all_rows[0].keys()); w.writeheader(); w.writerows(all_rows)
            return f"✅ 已合併 {len(all_rows)} 筆資料 → {out}"
        elif action == "to_table":
            with open(path, encoding="utf-8-sig", errors="replace") as f:
                rows = list(csv.DictReader(f))
            if not rows: return "⚠️ 無資料"
            cols = list(rows[0].keys())
            header = " | ".join(cols)
            sep = "-" * len(header)
            lines = [header, sep] + [" | ".join(str(r.get(c,"")) for c in cols) for r in rows[:20]]
            return "\n".join(lines)
    except Exception as e:
        return f"❌ 資料處理失敗：{e}"


def execute_wake_on_lan(mac, broadcast="255.255.255.255", port=9):
    try:
        import socket, struct
        mac_clean = mac.replace(":", "").replace("-", "")
        if len(mac_clean) != 12:
            return f"❌ MAC 位址格式錯誤：{mac}"
        mac_bytes = bytes.fromhex(mac_clean)
        magic = b"\xff" * 6 + mac_bytes * 16
        with socket.socket(socket.AF_INET, socket.SOCK_DGRAM) as sock:
            sock.setsockopt(socket.SOL_SOCKET, socket.SO_BROADCAST, 1)
            sock.sendto(magic, (broadcast, int(port)))
        return f"✅ WOL 魔法封包已送出 → {mac}（{broadcast}:{port}）"
    except Exception as e:
        return f"❌ WOL 失敗：{e}"


def execute_clipboard_history(action, index=0):
    global _clipboard_hist
    try:
        import pyperclip, threading, time
        if action == "start_watch":
            def _watch():
                last = ""
                while True:
                    try:
                        cur = pyperclip.paste()
                        if cur != last and cur:
                            _clipboard_hist.insert(0, cur)
                            if len(_clipboard_hist) > 50:
                                _clipboard_hist.pop()
                            last = cur
                    except Exception: pass
                    time.sleep(1)
            threading.Thread(target=_watch, daemon=True).start()
            return "✅ 剪貼簿歷史監控已啟動"
        elif action == "stop_watch":
            return "✅ 剪貼簿監控已停止（重啟才完全生效）"
        elif action == "list":
            if not _clipboard_hist:
                return "⚠️ 剪貼簿歷史為空（請先執行 start_watch）"
            lines = [f"[{i}] {item[:80]}" for i, item in enumerate(_clipboard_hist[:20])]
            return "📋 剪貼簿歷史：\n" + "\n".join(lines)
        elif action == "get":
            if index < len(_clipboard_hist):
                return f"📋 [{index}]：{_clipboard_hist[index]}"
            return f"⚠️ 索引 {index} 超出範圍（共 {len(_clipboard_hist)} 筆）"
        elif action == "set":
            if index < len(_clipboard_hist):
                import pyperclip
                pyperclip.copy(_clipboard_hist[index])
                return f"✅ 已復原剪貼簿 [{index}]：{_clipboard_hist[index][:80]}"
            return f"⚠️ 索引 {index} 超出範圍"
        elif action == "clear":
            _clipboard_hist.clear()
            return "✅ 剪貼簿歷史已清除"
    except Exception as e:
        return f"❌ 剪貼簿歷史失敗：{e}"


_file_watchers = {}
_pixel_watchers = {}
_mouse_recordings = {}
_screen_live_running = False

def execute_file_watcher(action, name="", path="", events="all", command="", notify=True, _bot_send=None, _chat_id=None):
    global _file_watchers
    try:
        import threading
        if action == "list":
            if not _file_watchers:
                return "⚠️ 無執行中的監聽器"
            return "📁 檔案監聽器：\n" + "\n".join(f"- {k}: {v['path']}" for k,v in _file_watchers.items())
        elif action == "stop":
            if name in _file_watchers:
                _file_watchers[name]["observer"].stop()
                del _file_watchers[name]
                return f"✅ 已停止監聽：{name}"
            return f"⚠️ 找不到：{name}"
        elif action == "start":
            from watchdog.observers import Observer
            from watchdog.events import FileSystemEventHandler
            class _Handler(FileSystemEventHandler):
                def _handle(self, event, etype):
                    if events != "all" and etype not in events: return
                    msg = f"📁 [{name}] {etype}：{event.src_path}"
                    if command:
                        import subprocess; subprocess.Popen(command.replace("{path}", event.src_path), shell=True)
                    if notify and _bot_send and _chat_id:
                        import asyncio
                        asyncio.run_coroutine_threadsafe(_bot_send(_chat_id, msg), asyncio.get_event_loop())
                def on_created(self, e): self._handle(e, "created")
                def on_modified(self, e): self._handle(e, "modified")
                def on_deleted(self, e): self._handle(e, "deleted")
            observer = Observer()
            observer.schedule(_Handler(), path, recursive=True)
            observer.start()
            _file_watchers[name] = {"path": path, "observer": observer}
            return f"✅ 已開始監聽：{name} → {path}（事件：{events}）"
    except Exception as e:
        return f"❌ 檔案監聽失敗：{e}"


def execute_pixel_watch(action, name="", x=0, y=0, command="", interval=1.0, tolerance=10, _bot_send=None, _chat_id=None):
    global _pixel_watchers
    try:
        import threading, time
        if action == "get":
            import pyautogui
            screenshot = pyautogui.screenshot()
            r, g, b = screenshot.getpixel((int(x), int(y)))[:3]
            return f"🎨 座標({x},{y}) 目前顏色：RGB({r},{g},{b}) #{r:02X}{g:02X}{b:02X}"
        elif action == "list":
            if not _pixel_watchers:
                return "⚠️ 無執行中的像素監控"
            return "🎨 像素監控：\n" + "\n".join(f"- {k}: ({v['x']},{v['y']})" for k,v in _pixel_watchers.items())
        elif action == "stop":
            if name in _pixel_watchers:
                _pixel_watchers[name]["running"] = False
                del _pixel_watchers[name]
                return f"✅ 已停止像素監控：{name}"
            return f"⚠️ 找不到：{name}"
        elif action == "start":
            import pyautogui
            screenshot = pyautogui.screenshot()
            r0, g0, b0 = screenshot.getpixel((int(x), int(y)))[:3]
            cfg = {"x": x, "y": y, "r": r0, "g": g0, "b": b0, "running": True}
            _pixel_watchers[name] = cfg
            def _watch():
                import pyautogui, subprocess, time as t
                while _pixel_watchers.get(name, {}).get("running"):
                    try:
                        sc = pyautogui.screenshot()
                        r, g, b = sc.getpixel((int(x), int(y)))[:3]
                        cfg = _pixel_watchers.get(name, {})
                        diff = abs(r-cfg["r"]) + abs(g-cfg["g"]) + abs(b-cfg["b"])
                        if diff > int(tolerance) * 3:
                            msg = f"🎨 [{name}] 像素({x},{y})顏色變化！#{r:02X}{g:02X}{b:02X}"
                            if command: subprocess.Popen(command, shell=True)
                            if _bot_send and _chat_id:
                                import asyncio
                                asyncio.run_coroutine_threadsafe(_bot_send(_chat_id, msg), asyncio.get_event_loop())
                            cfg["r"], cfg["g"], cfg["b"] = r, g, b
                    except Exception: pass
                    t.sleep(float(interval))
            threading.Thread(target=_watch, daemon=True).start()
            return f"✅ 像素監控已啟動：{name} ({x},{y}) 容差={tolerance}"
    except Exception as e:
        return f"❌ 像素監控失敗：{e}"


def execute_object_detect(target, action="find", region=""):
    try:
        import pyautogui, anthropic, base64, io as _io, json, re, os
        reg = None
        if region:
            parts = [int(v) for v in region.split(",")]
            if len(parts) == 4: reg = tuple(parts)
        screenshot = pyautogui.screenshot(region=reg)
        buf = _io.BytesIO(); screenshot.save(buf, format="PNG")
        img_b64 = base64.standard_b64encode(buf.getvalue()).decode()
        offset_x = reg[0] if reg else 0
        offset_y = reg[1] if reg else 0
        _client = anthropic.Anthropic(api_key=os.getenv("ANTHROPIC_API_KEY"))
        resp = _client.messages.create(
            model="claude-sonnet-4-6", max_tokens=256,
            messages=[{"role": "user", "content": [
                {"type": "image", "source": {"type": "base64", "media_type": "image/png", "data": img_b64}},
                {"type": "text", "text": f"在圖片中找到「{target}」的位置。回答 JSON：{{\"found\": true/false, \"x\": 中心X座標, \"y\": 中心Y座標, \"description\": \"說明\"}}"}
            ]}]
        )
        m = re.search(r'\{.*\}', resp.content[0].text, re.DOTALL)
        if not m: return f"⚠️ AI 無法解析回應"
        result = json.loads(m.group())
        if not result.get("found"): return f"⚠️ 未找到：{target}"
        ax = result["x"] + offset_x
        ay = result["y"] + offset_y
        if action == "click": pyautogui.click(ax, ay); return f"✅ 已點擊「{target}」({ax},{ay})"
        elif action == "double_click": pyautogui.doubleClick(ax, ay); return f"✅ 已雙擊「{target}」({ax},{ay})"
        return f"✅ 找到「{target}」：座標({ax},{ay}) — {result.get('description','')}"
    except Exception as e:
        return f"❌ 物件偵測失敗：{e}"


def execute_mouse_record(action, name="", duration=10.0, repeat=1, speed=1.0):
    try:
        import json, time
        from pathlib import Path as _Path
        record_file = _Path.home() / ".claude_mouse_macros.json"
        store = json.loads(record_file.read_text()) if record_file.exists() else {}
        if action == "list":
            if not store: return "⚠️ 無已儲存的滑鼠巨集"
            return "🖱️ 滑鼠巨集：\n" + "\n".join(f"- {k}（{len(v)} 個事件）" for k,v in store.items())
        elif action == "delete":
            if name in store:
                del store[name]; record_file.write_text(json.dumps(store))
                return f"✅ 已刪除：{name}"
            return f"⚠️ 找不到：{name}"
        elif action == "start":
            from pynput import mouse, keyboard as kb
            events = []
            start_time = time.time()
            def on_move(x, y): events.append({"t": time.time()-start_time, "type":"move", "x":x, "y":y})
            def on_click(x, y, btn, pressed): events.append({"t": time.time()-start_time, "type":"click", "x":x, "y":y, "btn":str(btn), "pressed":pressed})
            def on_scroll(x, y, dx, dy): events.append({"t": time.time()-start_time, "type":"scroll", "x":x, "y":y, "dx":dx, "dy":dy})
            m_listener = mouse.Listener(on_move=on_move, on_click=on_click, on_scroll=on_scroll)
            m_listener.start()
            time.sleep(float(duration))
            m_listener.stop()
            store[name] = events
            record_file.write_text(json.dumps(store))
            return f"✅ 滑鼠巨集 '{name}' 錄製完成（{len(events)} 個事件，{duration}s）"
        elif action == "play":
            if name not in store: return f"⚠️ 找不到：{name}"
            from pynput import mouse as m
            controller = m.Controller()
            events = store[name]
            spd = float(speed)
            for _ in range(int(repeat)):
                prev_t = 0.0
                for e in events:
                    delay = (e["t"] - prev_t) / spd
                    if delay > 0: time.sleep(min(delay, 0.5))
                    prev_t = e["t"]
                    if e["type"] == "move":
                        controller.position = (e["x"], e["y"])
                    elif e["type"] == "click":
                        btn = m.Button.left if "left" in e["btn"] else m.Button.right
                        if e["pressed"]: controller.press(btn)
                        else: controller.release(btn)
                    elif e["type"] == "scroll":
                        controller.scroll(e["dx"], e["dy"])
            return f"✅ 滑鼠巨集 '{name}' 回放 {repeat} 次完成"
    except Exception as e:
        return f"❌ 滑鼠巨集失敗：{e}"


def execute_adb(action, x=0, y=0, x2=0, y2=0, text="", path="", remote="", package="", command="", device=""):
    try:
        import subprocess
        prefix = ["adb"]
        if device: prefix += ["-s", device]
        if action == "devices":
            r = subprocess.run(["adb", "devices", "-l"], capture_output=True, text=True)
            return f"📱 ADB 裝置：\n{r.stdout.strip()}"
        elif action == "screenshot":
            out = path or str(Path("C:/Users/blue_/Desktop/測試檔案") / f"adb_{datetime.datetime.now().strftime('%H%M%S')}.png")
            subprocess.run(prefix + ["shell", "screencap", "-p", "/sdcard/screen.png"], capture_output=True)
            subprocess.run(prefix + ["pull", "/sdcard/screen.png", out], capture_output=True)
            return f"✅ 手機截圖已存：{out}"
        elif action == "tap":
            subprocess.run(prefix + ["shell", "input", "tap", str(x), str(y)], capture_output=True)
            return f"✅ 已點擊手機 ({x},{y})"
        elif action == "swipe":
            subprocess.run(prefix + ["shell", "input", "swipe", str(x), str(y), str(x2), str(y2), "300"], capture_output=True)
            return f"✅ 已滑動 ({x},{y})→({x2},{y2})"
        elif action == "type":
            t = text.replace(" ", "%s")
            subprocess.run(prefix + ["shell", "input", "text", t], capture_output=True)
            return f"✅ 已輸入文字：{text}"
        elif action == "key":
            subprocess.run(prefix + ["shell", "input", "keyevent", text], capture_output=True)
            return f"✅ 已按鍵：{text}"
        elif action == "install":
            r = subprocess.run(prefix + ["install", "-r", path], capture_output=True, text=True)
            return f"✅ 已安裝：{path}\n{r.stdout.strip()}" if r.returncode == 0 else f"❌ 安裝失敗：{r.stderr}"
        elif action == "push":
            r = subprocess.run(prefix + ["push", path, remote or "/sdcard/"], capture_output=True, text=True)
            return f"✅ 已上傳：{path} → {remote}" if r.returncode == 0 else f"❌ 失敗：{r.stderr}"
        elif action == "pull":
            out = path or str(Path("C:/Users/blue_/Desktop/測試檔案") / Path(remote).name)
            r = subprocess.run(prefix + ["pull", remote, out], capture_output=True, text=True)
            return f"✅ 已下載：{remote} → {out}" if r.returncode == 0 else f"❌ 失敗：{r.stderr}"
        elif action == "shell":
            r = subprocess.run(prefix + ["shell", command], capture_output=True, text=True)
            return f"📱 ADB Shell：\n{r.stdout.strip()}"
        elif action == "start_app":
            r = subprocess.run(prefix + ["shell", "monkey", "-p", package, "-c", "android.intent.category.LAUNCHER", "1"], capture_output=True, text=True)
            return f"✅ 已啟動 App：{package}"
        elif action == "stop_app":
            subprocess.run(prefix + ["shell", "am", "force-stop", package], capture_output=True)
            return f"✅ 已停止 App：{package}"
    except Exception as e:
        return f"❌ ADB 操作失敗：{e}（請確認已安裝 ADB 並開啟 USB 偵錯）"


def execute_wifi_hotspot(action, ssid="", password=""):
    try:
        import subprocess
        if action == "status":
            r = subprocess.run(["powershell", "-Command",
                "Get-NetConnectionProfile | Select-Object Name,NetworkCategory | Format-Table; netsh wlan show hostednetwork"],
                capture_output=True, text=True, encoding="utf-8", errors="replace")
            return f"📡 熱點狀態：\n{r.stdout.strip()}"
        elif action == "set":
            r = subprocess.run(["netsh", "wlan", "set", "hostednetwork",
                f"mode=allow", f"ssid={ssid}", f"key={password}"],
                capture_output=True, text=True, encoding="utf-8", errors="replace")
            return f"✅ 熱點設定完成：SSID={ssid}\n{r.stdout.strip()}"
        elif action == "start":
            r = subprocess.run(["netsh", "wlan", "start", "hostednetwork"],
                capture_output=True, text=True, encoding="utf-8", errors="replace")
            return f"✅ 熱點已啟動\n{r.stdout.strip()}"
        elif action == "stop":
            r = subprocess.run(["netsh", "wlan", "stop", "hostednetwork"],
                capture_output=True, text=True, encoding="utf-8", errors="replace")
            return f"✅ 熱點已停止\n{r.stdout.strip()}"
    except Exception as e:
        return f"❌ WiFi 熱點失敗：{e}"


def execute_onedrive(action, path="", remote=""):
    try:
        import subprocess, os
        onedrive_path = os.path.expandvars(r"%USERPROFILE%\OneDrive")
        if not Path(onedrive_path).exists():
            onedrive_path = os.path.expandvars(r"%OneDrive%") or str(Path.home() / "OneDrive")
        if action == "status":
            r = subprocess.run(["powershell", "-Command",
                "Get-Process OneDrive -ErrorAction SilentlyContinue | Select-Object Name,CPU,WorkingSet"],
                capture_output=True, text=True, encoding="utf-8", errors="replace")
            size = sum(f.stat().st_size for f in Path(onedrive_path).rglob("*") if f.is_file()) / 1024 / 1024 / 1024
            return f"☁️ OneDrive 路徑：{onedrive_path}\n使用空間：{size:.2f} GB\n程序狀態：\n{r.stdout.strip()}"
        elif action == "list":
            target = Path(onedrive_path) / (remote or "")
            if not target.exists(): return f"⚠️ 路徑不存在：{target}"
            items = list(target.iterdir())
            lines = [f"{'📁' if p.is_dir() else '📄'} {p.name}" for p in sorted(items)]
            return f"☁️ OneDrive/{remote or ''}：\n" + "\n".join(lines[:30])
        elif action == "upload":
            import shutil
            dest = Path(onedrive_path) / (remote or Path(path).name)
            shutil.copy2(path, dest)
            return f"✅ 已上傳至 OneDrive：{dest}"
        elif action == "download":
            import shutil
            src = Path(onedrive_path) / remote
            if not src.exists(): return f"⚠️ 找不到：{src}"
            dest = path or str(Path("C:/Users/blue_/Desktop/測試檔案") / src.name)
            shutil.copy2(src, dest)
            return f"✅ 已從 OneDrive 下載：{dest}"
        elif action == "sync":
            subprocess.Popen(["powershell", "-Command",
                'Start-Process "$env:LOCALAPPDATA\\Microsoft\\OneDrive\\OneDrive.exe"'])
            return "✅ OneDrive 同步已觸發"
        elif action == "open":
            import os; os.startfile(onedrive_path)
            return f"✅ 已開啟 OneDrive 資料夾：{onedrive_path}"
    except Exception as e:
        return f"❌ OneDrive 操作失敗：{e}"


def execute_ftp(action, host="", user="", password="", local="", remote="", port=21):
    try:
        from ftplib import FTP
        ftp = FTP()
        ftp.connect(host, int(port), timeout=30)
        ftp.login(user, password)
        if action == "list":
            items = ftp.nlst(remote or ".")
            ftp.quit()
            return f"📂 FTP {host}{remote or '/'}：\n" + "\n".join(items[:50])
        elif action == "upload":
            with open(local, "rb") as f:
                ftp.storbinary(f"STOR {remote or Path(local).name}", f)
            ftp.quit()
            return f"✅ 已上傳：{local} → {host}/{remote}"
        elif action == "download":
            out = local or str(Path("C:/Users/blue_/Desktop/測試檔案") / Path(remote).name)
            with open(out, "wb") as f:
                ftp.retrbinary(f"RETR {remote}", f.write)
            ftp.quit()
            return f"✅ 已下載：{host}/{remote} → {out}"
        elif action == "delete":
            ftp.delete(remote); ftp.quit()
            return f"✅ 已刪除：{remote}"
        elif action == "mkdir":
            ftp.mkd(remote); ftp.quit()
            return f"✅ 已建立目錄：{remote}"
        elif action == "rename":
            ftp.rename(remote, local); ftp.quit()
            return f"✅ 已重新命名：{remote} → {local}"
    except Exception as e:
        return f"❌ FTP 操作失敗：{e}"


def execute_wsl(action, distro="", command=""):
    try:
        import subprocess
        if action == "list":
            r = subprocess.run(["wsl", "--list", "--verbose"],
                capture_output=True, text=True, encoding="utf-16-le", errors="replace")
            return f"🐧 WSL 發行版：\n{r.stdout.strip()}"
        elif action == "status":
            r = subprocess.run(["wsl", "--status"],
                capture_output=True, text=True, encoding="utf-16-le", errors="replace")
            return f"🐧 WSL 狀態：\n{r.stdout.strip()}"
        elif action == "run":
            cmd = ["wsl"]
            if distro: cmd += ["-d", distro]
            cmd += ["--", "bash", "-c", command]
            r = subprocess.run(cmd, capture_output=True, text=True, encoding="utf-8", errors="replace")
            return f"🐧 WSL 輸出：\n{r.stdout.strip()}" + (f"\n錯誤：{r.stderr.strip()}" if r.stderr.strip() else "")
        elif action == "start":
            subprocess.Popen(["wsl"] + (["-d", distro] if distro else []))
            return f"✅ WSL 已啟動：{distro or '預設'}"
        elif action == "stop":
            cmd = ["wsl", "--terminate", distro] if distro else ["wsl", "--shutdown"]
            subprocess.run(cmd, capture_output=True)
            return f"✅ WSL 已停止：{distro or '全部'}"
        elif action == "install":
            r = subprocess.run(["wsl", "--install", "-d", distro],
                capture_output=True, text=True)
            return f"✅ 正在安裝 {distro}（需要網路）"
    except Exception as e:
        return f"❌ WSL 操作失敗：{e}（請確認已啟用 WSL）"


def execute_hyperv(action, name="", snapshot=""):
    try:
        import subprocess
        def ps(cmd):
            r = subprocess.run(["powershell", "-Command", cmd],
                capture_output=True, text=True, encoding="utf-8", errors="replace")
            return r.stdout.strip(), r.returncode
        if action == "list":
            out, _ = ps("Get-VM | Select-Object Name,State,CPUUsage,MemoryAssigned | Format-Table -AutoSize")
            return f"💻 虛擬機清單：\n{out}" if out else "⚠️ 未找到虛擬機（請確認已啟用 Hyper-V）"
        elif action == "status":
            out, _ = ps(f"Get-VM -Name '{name}' | Select-Object * | Format-List")
            return f"💻 {name} 狀態：\n{out}"
        elif action == "start":
            out, rc = ps(f"Start-VM -Name '{name}'")
            return f"✅ 已啟動：{name}" if rc == 0 else f"❌ 啟動失敗：{out}"
        elif action == "stop":
            out, rc = ps(f"Stop-VM -Name '{name}' -Force")
            return f"✅ 已停止：{name}" if rc == 0 else f"❌ 停止失敗：{out}"
        elif action == "pause":
            ps(f"Suspend-VM -Name '{name}'")
            return f"✅ 已暫停：{name}"
        elif action == "resume":
            ps(f"Resume-VM -Name '{name}'")
            return f"✅ 已繼續：{name}"
        elif action == "snapshot":
            sname = snapshot or datetime.datetime.now().strftime("snap_%Y%m%d_%H%M%S")
            out, rc = ps(f"Checkpoint-VM -Name '{name}' -SnapshotName '{sname}'")
            return f"✅ 快照已建立：{sname}" if rc == 0 else f"❌ 失敗：{out}"
        elif action == "restore":
            out, rc = ps(f"Restore-VMSnapshot -Name '{name}' -VMName '{name}' -Confirm:$false")
            return f"✅ 已還原快照：{snapshot}" if rc == 0 else f"❌ 失敗：{out}"
        elif action == "delete_snapshot":
            ps(f"Remove-VMSnapshot -VMName '{name}' -Name '{snapshot}' -Confirm:$false")
            return f"✅ 已刪除快照：{snapshot}"
    except Exception as e:
        return f"❌ Hyper-V 操作失敗：{e}"


def execute_file_diff(file1, file2, output="", mode="unified"):
    try:
        import difflib
        text1 = Path(file1).read_text(encoding="utf-8", errors="replace").splitlines(keepends=True)
        text2 = Path(file2).read_text(encoding="utf-8", errors="replace").splitlines(keepends=True)
        if mode == "unified":
            diff = list(difflib.unified_diff(text1, text2, fromfile=file1, tofile=file2))
        elif mode == "context":
            diff = list(difflib.context_diff(text1, text2, fromfile=file1, tofile=file2))
        else:
            diff = [f"{'+ ' if a != b else '  '}{a}" for a, b in zip(text1, text2)]
        if not diff:
            return "✅ 兩個檔案內容完全相同"
        result = "".join(diff)
        if output:
            Path(output).write_text(result, encoding="utf-8")
            return f"✅ diff 已儲存：{output}\n（共 {len(diff)} 行差異）"
        return f"📄 檔案差異（{len(diff)} 行）：\n{result[:2000]}"
    except Exception as e:
        return f"❌ 檔案比較失敗：{e}"


def execute_screen_live(action, fps=0.5, duration=60.0, quality=50, _bot_send=None, _chat_id=None):
    global _screen_live_running
    try:
        import threading
        if action == "stop":
            _screen_live_running = False
            return "✅ 螢幕串流已停止"
        elif action == "start":
            if _screen_live_running:
                return "⚠️ 螢幕串流已在執行中"
            _screen_live_running = True
            def _stream():
                global _screen_live_running
                import pyautogui, io as _io, time as t, asyncio
                interval = 1.0 / max(float(fps), 0.1)
                end = t.time() + float(duration)
                loop = asyncio.get_event_loop()
                count = 0
                while _screen_live_running and t.time() < end:
                    try:
                        screenshot = pyautogui.screenshot()
                        buf = _io.BytesIO()
                        screenshot.save(buf, format="JPEG", quality=int(quality))
                        buf.seek(0)
                        if _bot_send and _chat_id:
                            import telegram
                            asyncio.run_coroutine_threadsafe(
                                _bot_send(_chat_id, photo=buf),
                                loop)
                        count += 1
                    except Exception:
                        pass
                    t.sleep(interval)
                _screen_live_running = False
            threading.Thread(target=_stream, daemon=True).start()
            return f"✅ 螢幕串流已啟動（{fps} FPS，{duration}s，畫質 {quality}）"
    except Exception as e:
        return f"❌ 螢幕串流失敗：{e}"


def execute_download_file(url, save_path=""):
    try:
        import requests
        if not save_path:
            fname = url.split("/")[-1].split("?")[0] or "download"
            save_path = str(Path("C:/Users/blue_/Desktop/測試檔案") / fname)
        r = requests.get(url, stream=True, timeout=60)
        r.raise_for_status()
        with open(save_path, "wb") as f:
            for chunk in r.iter_content(chunk_size=8192):
                f.write(chunk)
        return f"✅ 已下載：{save_path}"
    except Exception as e:
        return f"❌ 下載失敗：{e}"


def execute_wake_listen(keyword="小牛馬", duration=5):
    try:
        import sounddevice as sd, soundfile as sf, speech_recognition as sr, tempfile, time
        sample_rate = 16000
        deadline = time.time() + 60
        while time.time() < deadline:
            recording = sd.rec(int(int(duration) * sample_rate), samplerate=sample_rate, channels=1, dtype="int16")
            sd.wait()
            with tempfile.NamedTemporaryFile(suffix=".wav", delete=False) as f:
                tmp_path = f.name
            sf.write(tmp_path, recording, sample_rate)
            recognizer = sr.Recognizer()
            try:
                with sr.AudioFile(tmp_path) as source:
                    audio = recognizer.record(source)
                text = recognizer.recognize_google(audio, language="zh-TW")
                Path(tmp_path).unlink(missing_ok=True)
                if keyword in text:
                    return f"✅ 偵測到喚醒詞！辨識到：{text}"
            except Exception:
                Path(tmp_path).unlink(missing_ok=True)
        return f"⏳ 60 秒內未偵測到喚醒詞「{keyword}」"
    except Exception as e:
        return f"❌ 語音監聽失敗：{e}"


def execute_right_menu(x, y, item=""):
    try:
        import pyautogui, time
        pyautogui.rightClick(int(x), int(y))
        time.sleep(0.3)
        if item:
            pyautogui.write(item, interval=0.05)
            time.sleep(0.2)
            pyautogui.press("enter")
            return f"✅ 已右鍵點擊並選擇：{item}"
        return f"✅ 已右鍵點擊 ({x},{y})"
    except Exception as e:
        return f"❌ 右鍵選單失敗：{e}"


def execute_disk_clean(action="list"):
    try:
        import tempfile, shutil
        tmp = Path(tempfile.gettempdir())
        if action == "list":
            files = list(tmp.rglob("*"))
            total = sum(f.stat().st_size for f in files if f.is_file())
            return f"🗑️ 暫存資料夾：{tmp}\n檔案數：{len(files)}\n佔用：{total/1024/1024:.1f} MB"
        elif action == "clean":
            count = 0
            for f in tmp.iterdir():
                try:
                    if f.is_file(): f.unlink(); count += 1
                    elif f.is_dir(): shutil.rmtree(f, ignore_errors=True); count += 1
                except Exception: pass
            return f"✅ 已清理 {count} 個暫存項目"
    except Exception as e:
        return f"❌ 磁碟清理失敗：{e}"


def execute_usb_list():
    try:
        import subprocess
        r = subprocess.run(["powershell", "-Command",
            "Get-PnpDevice | Where-Object {$_.Class -eq 'USB' -and $_.Status -eq 'OK'} | Select-Object FriendlyName,InstanceId | Format-Table -AutoSize"],
            capture_output=True, text=True, encoding="utf-8", errors="replace")
        return f"🔌 USB 裝置：\n{r.stdout.strip()[:2000]}" if r.stdout.strip() else "⚠️ 無 USB 裝置"
    except Exception as e:
        return f"❌ USB 查詢失敗：{e}"


def execute_rdp_connect(host, user="", width=1280, height=720):
    try:
        import subprocess
        args = ["/v:" + host, f"/w:{width}", f"/h:{height}"]
        if user:
            args.append(f"/u:{user}")
        subprocess.Popen(["mstsc"] + args)
        return f"✅ 正在連線 RDP：{host}"
    except Exception as e:
        return f"❌ RDP 連線失敗：{e}"


def execute_chrome_bookmarks():
    try:
        import json
        bookmark_path = Path.home() / "AppData/Local/Google/Chrome/User Data/Default/Bookmarks"
        if not bookmark_path.exists():
            return "❌ 找不到 Chrome 書籤（Chrome 未安裝或路徑不同）"
        data = json.loads(bookmark_path.read_text(encoding="utf-8"))
        lines = []
        def _collect(node, indent=0):
            if node.get("type") == "url":
                lines.append("  " * indent + f"🔗 {node['name']}  {node['url']}")
            elif node.get("type") == "folder":
                lines.append("  " * indent + f"📁 {node['name']}")
                for child in node.get("children", []):
                    _collect(child, indent + 1)
        for root in data["roots"].values():
            _collect(root)
        return "📚 Chrome 書籤：\n" + "\n".join(lines[:80])
    except Exception as e:
        return f"❌ 讀取書籤失敗：{e}"


def execute_net_share(action, share_path="", drive="Z:", user="", password=""):
    try:
        import subprocess
        if action == "list":
            r = subprocess.run(["net", "use"], capture_output=True, text=True, encoding="cp950", errors="replace")
            return f"🌐 網路磁碟機：\n{r.stdout.strip()}"
        elif action == "connect":
            args = ["net", "use", drive, share_path]
            if user: args += [f"/user:{user}", password]
            r = subprocess.run(args, capture_output=True, text=True, encoding="cp950", errors="replace")
            return r.stdout.strip() or f"✅ 已連線 {share_path} → {drive}"
        elif action == "disconnect":
            r = subprocess.run(["net", "use", drive, "/delete"], capture_output=True, text=True, encoding="cp950", errors="replace")
            return r.stdout.strip() or f"✅ 已中斷 {drive}"
    except Exception as e:
        return f"❌ 網路芳鄰失敗：{e}"


def execute_font_list(keyword=""):
    try:
        import subprocess
        r = subprocess.run(["powershell", "-Command",
            "[System.Reflection.Assembly]::LoadWithPartialName('System.Drawing') | Out-Null; [System.Drawing.FontFamily]::Families | Select-Object -ExpandProperty Name"],
            capture_output=True, text=True, encoding="utf-8", errors="replace")
        fonts = r.stdout.strip().splitlines()
        if keyword:
            fonts = [f for f in fonts if keyword.lower() in f.lower()]
        result = "\n".join(fonts[:50])
        suffix = f"\n...（共 {len(fonts)} 個字型）" if len(fonts) > 50 else f"\n共 {len(fonts)} 個字型"
        return f"🔤 字型清單：\n{result}{suffix}"
    except Exception as e:
        return f"❌ 字型查詢失敗：{e}"


def execute_wait_seconds(seconds):
    try:
        import time
        s = float(seconds)
        time.sleep(min(s, 60))
        return f"✅ 已等待 {s} 秒"
    except Exception as e:
        return f"❌ 等待失敗：{e}"


def execute_firewall(action, name="", port=None, protocol="TCP", direction="Inbound"):
    try:
        import subprocess
        if action == "status":
            r = subprocess.run(["powershell", "-Command",
                "Get-NetFirewallProfile | Select-Object Name,Enabled | Format-Table -AutoSize"],
                capture_output=True, text=True, encoding="utf-8", errors="replace")
            return f"🔥 防火牆狀態：\n{r.stdout.strip()}"
        elif action == "list":
            r = subprocess.run(["powershell", "-Command",
                f"Get-NetFirewallRule | Where-Object {{$_.Enabled -eq 'True'}} | Select-Object DisplayName,Direction,Action | Sort-Object Direction | Format-Table -AutoSize"],
                capture_output=True, text=True, encoding="utf-8", errors="replace")
            return f"🔥 防火牆規則：\n{r.stdout.strip()[:2000]}"
        elif action == "add":
            r = subprocess.run(["powershell", "-Command",
                f"New-NetFirewallRule -DisplayName '{name}' -Direction {direction} -Protocol {protocol} -LocalPort {port} -Action Allow -Enabled True"],
                capture_output=True, text=True)
            return f"✅ 已新增防火牆規則：{name} ({direction} {protocol}:{port})" if r.returncode == 0 else f"❌ 失敗：{r.stderr.strip()}"
        elif action == "remove":
            r = subprocess.run(["powershell", "-Command",
                f"Remove-NetFirewallRule -DisplayName '{name}' -Confirm:$false"],
                capture_output=True, text=True)
            return f"✅ 已移除規則：{name}" if r.returncode == 0 else f"❌ 失敗：{r.stderr.strip()}"
        elif action == "enable":
            subprocess.run(["powershell", "-Command",
                "Set-NetFirewallProfile -Profile Domain,Public,Private -Enabled True"], capture_output=True)
            return "✅ 防火牆已啟用"
        elif action == "disable":
            subprocess.run(["powershell", "-Command",
                "Set-NetFirewallProfile -Profile Domain,Public,Private -Enabled False"], capture_output=True)
            return "✅ 防火牆已停用"
    except Exception as e:
        return f"❌ 防火牆操作失敗：{e}"


def execute_process_mgr(action, name="", pid=None, level="normal"):
    try:
        import psutil
        priority_map = {
            "realtime": psutil.REALTIME_PRIORITY_CLASS,
            "high": psutil.HIGH_PRIORITY_CLASS,
            "above_normal": psutil.ABOVE_NORMAL_PRIORITY_CLASS,
            "normal": psutil.NORMAL_PRIORITY_CLASS,
            "below_normal": psutil.BELOW_NORMAL_PRIORITY_CLASS,
            "idle": psutil.IDLE_PRIORITY_CLASS,
        }
        if action == "list":
            procs = sorted(psutil.process_iter(["pid","name","cpu_percent","memory_info"]), key=lambda p: p.info["cpu_percent"] or 0, reverse=True)
            lines = [f"{'PID':>6} {'CPU%':>6} {'MEM MB':>8}  名稱"]
            for p in procs[:25]:
                mem = (p.info["memory_info"].rss // 1024 // 1024) if p.info["memory_info"] else 0
                lines.append(f"{p.info['pid']:>6} {p.info['cpu_percent'] or 0:>6.1f} {mem:>8}  {p.info['name']}")
            return "\n".join(lines)
        elif action == "search":
            found = [p for p in psutil.process_iter(["pid","name","cpu_percent","memory_info"]) if name.lower() in p.info["name"].lower()]
            if not found:
                return f"⚠️ 找不到程序：{name}"
            lines = [f"PID:{p.info['pid']} CPU:{p.info['cpu_percent']}% 記憶體:{p.info['memory_info'].rss//1024//1024}MB {p.info['name']}" for p in found]
            return "\n".join(lines)
        elif action == "kill":
            targets = []
            if pid:
                targets = [psutil.Process(int(pid))]
            else:
                targets = [p for p in psutil.process_iter(["pid","name"]) if name.lower() in p.info["name"].lower()]
            if not targets:
                return f"⚠️ 找不到程序：{name}"
            for p in targets:
                p.kill()
            return f"✅ 已終止 {len(targets)} 個程序：{name or pid}"
        elif action == "priority":
            p = psutil.Process(int(pid)) if pid else next((x for x in psutil.process_iter(["pid","name"]) if name.lower() in x.info["name"].lower()), None)
            if not p:
                return f"⚠️ 找不到程序：{name}"
            p.nice(priority_map.get(level, psutil.NORMAL_PRIORITY_CLASS))
            return f"✅ 已設定 PID {p.pid} 優先權為：{level}"
    except Exception as e:
        return f"❌ 程序管理失敗：{e}"


def execute_power_plan(action, plan="balanced"):
    try:
        import subprocess
        plan_guids = {
            "balanced": "381b4222-f694-41f0-9685-ff5bb260df2e",
            "high_performance": "8c5e7fda-e8bf-4a96-9a85-a6e23a8c635c",
            "power_saver": "a1841308-3541-4fab-bc81-f71556f20b4a",
        }
        if action == "list":
            r = subprocess.run(["powercfg", "/list"], capture_output=True, text=True, encoding="utf-8", errors="replace")
            return f"⚡ 電源計畫：\n{r.stdout.strip()}"
        elif action == "get":
            r = subprocess.run(["powercfg", "/getactivescheme"], capture_output=True, text=True, encoding="utf-8", errors="replace")
            return f"⚡ 目前電源計畫：\n{r.stdout.strip()}"
        elif action == "set":
            guid = plan_guids.get(plan)
            if not guid:
                return f"⚠️ 未知計畫：{plan}"
            subprocess.run(["powercfg", "/setactive", guid], capture_output=True)
            return f"✅ 電源計畫已切換為：{plan}"
    except Exception as e:
        return f"❌ 電源計畫失敗：{e}"


def execute_event_log(log="System", level="Error", count=10):
    try:
        import win32evtlog, win32evtlogutil, win32con
        level_map = {"Error": 1, "Warning": 2, "Information": 4, "All": 7}
        event_type = level_map.get(level, 1)
        hand = win32evtlog.OpenEventLog(None, log)
        flags = win32evtlog.EVENTLOG_BACKWARDS_READ | win32evtlog.EVENTLOG_SEQUENTIAL_READ
        events = []
        while len(events) < int(count):
            batch = win32evtlog.ReadEventLog(hand, flags, 0)
            if not batch:
                break
            for e in batch:
                if level == "All" or e.EventType == event_type:
                    try:
                        msg = win32evtlogutil.SafeFormatMessage(e, log)[:100]
                    except Exception:
                        msg = "(無法讀取訊息)"
                    events.append(f"[{e.TimeGenerated.Format()}] {e.SourceName}: {msg}")
                if len(events) >= int(count):
                    break
        win32evtlog.CloseEventLog(hand)
        if not events:
            return f"✅ {log} 中沒有 {level} 等級事件"
        return f"📋 {log} 事件記錄（{level}）：\n" + "\n".join(events)
    except Exception as e:
        return f"❌ 事件記錄讀取失敗：{e}"


def execute_datetime_config(action, timezone="", datetime_str=""):
    try:
        import subprocess
        if action == "get":
            r = subprocess.run(["powershell", "-Command",
                "Get-Date | Format-List; (Get-TimeZone).DisplayName"],
                capture_output=True, text=True, encoding="utf-8", errors="replace")
            return f"🕐 系統時間：\n{r.stdout.strip()}"
        elif action == "sync":
            subprocess.run(["powershell", "-Command",
                "Start-Service w32tm -ErrorAction SilentlyContinue; w32tm /resync /force"],
                capture_output=True)
            return "✅ 已同步網路時間"
        elif action == "set_timezone":
            r = subprocess.run(["powershell", "-Command",
                f"Set-TimeZone -Id '{timezone}'"],
                capture_output=True, text=True)
            return f"✅ 時區已設定為：{timezone}" if r.returncode == 0 else f"❌ 設定失敗：{r.stderr.strip()}"
        elif action == "set_time":
            r = subprocess.run(["powershell", "-Command",
                f"Set-Date -Date '{datetime_str}'"],
                capture_output=True, text=True)
            return f"✅ 系統時間已設定為：{datetime_str}" if r.returncode == 0 else f"❌ 設定失敗：{r.stderr.strip()}"
    except Exception as e:
        return f"❌ 時間設定失敗：{e}"


def execute_ui_auto(action, window="", control="", text=""):
    try:
        from pywinauto import Desktop, Application
        desktop = Desktop(backend="uia")
        if action == "get_windows":
            wins = [str(w.window_text()) for w in desktop.windows() if w.window_text()]
            return "🪟 所有視窗：\n" + "\n".join(f"- {w}" for w in wins[:30])
        win = None
        for w in desktop.windows():
            if window.lower() in w.window_text().lower():
                win = w
                break
        if not win:
            return f"⚠️ 找不到視窗：{window}"
        if action == "find":
            ctrls = win.descendants()
            info = [f"[{c.control_type()}] {c.window_text()}" for c in ctrls if c.window_text()][:30]
            return f"🔍 視窗 '{window}' 的控制項：\n" + "\n".join(info)
        elif action == "read":
            texts = [c.window_text() for c in win.descendants() if c.window_text()]
            return f"📖 '{window}' 內容：\n" + "\n".join(texts[:50])
        elif action == "click":
            for c in win.descendants():
                if control.lower() in c.window_text().lower():
                    c.click_input()
                    return f"✅ 已點擊：{c.window_text()}"
            return f"⚠️ 找不到控制項：{control}"
        elif action == "type":
            for c in win.descendants():
                if control.lower() in c.window_text().lower() or c.control_type() in ("Edit","Document"):
                    c.type_keys(text)
                    return f"✅ 已輸入文字到：{c.window_text() or c.control_type()}"
            return f"⚠️ 找不到輸入框：{control}"
    except Exception as e:
        return f"❌ UI 自動化失敗：{e}"


_macro_recordings = {}
_macro_store = {}

def execute_macro(action, name="", repeat=1, duration=10.0):
    try:
        import keyboard, json, time
        global _macro_recordings, _macro_store
        macro_file = Path.home() / ".claude_macros.json"
        if macro_file.exists():
            _macro_store = json.loads(macro_file.read_text())

        if action == "record_start":
            events = []
            keyboard.start_recording()
            time.sleep(float(duration))
            recorded = keyboard.stop_recording()
            _macro_store[name] = [{"type": e.event_type, "name": e.name, "time": e.time} for e in recorded]
            macro_file.write_text(json.dumps(_macro_store))
            return f"✅ 巨集 '{name}' 錄製完成（{len(recorded)} 個事件，{duration}s）"
        elif action == "record_stop":
            return "⚠️ 請使用 record_start 並設定 duration 秒數"
        elif action == "play":
            if name not in _macro_store:
                return f"⚠️ 找不到巨集：{name}"
            events = _macro_store[name]
            for _ in range(int(repeat)):
                prev_time = events[0]["time"] if events else 0
                for e in events:
                    delay = max(0, e["time"] - prev_time)
                    time.sleep(min(delay, 0.5))
                    prev_time = e["time"]
                    if e["type"] == "down":
                        keyboard.press(e["name"])
                    elif e["type"] == "up":
                        keyboard.release(e["name"])
            return f"✅ 巨集 '{name}' 已回放 {repeat} 次"
        elif action == "list":
            if not _macro_store:
                return "⚠️ 無已儲存巨集"
            return "📋 已儲存巨集：\n" + "\n".join(f"- {k}（{len(v)} 事件）" for k,v in _macro_store.items())
        elif action == "delete":
            if name in _macro_store:
                del _macro_store[name]
                macro_file.write_text(json.dumps(_macro_store))
                return f"✅ 已刪除巨集：{name}"
            return f"⚠️ 找不到巨集：{name}"
    except Exception as e:
        return f"❌ 巨集操作失敗：{e}"


def execute_color_pick(action, x=0, y=0, region_w=100, region_h=100):
    try:
        import pyautogui
        from PIL import Image
        if action == "pick":
            screenshot = pyautogui.screenshot()
            pixel = screenshot.getpixel((int(x), int(y)))
            r, g, b = pixel[:3]
            hex_color = f"#{r:02X}{g:02X}{b:02X}"
            return f"🎨 座標 ({x},{y}) 的顏色：\nRGB: ({r}, {g}, {b})\nHEX: {hex_color}"
        elif action == "dominant":
            screenshot = pyautogui.screenshot(region=(int(x), int(y), int(region_w), int(region_h)))
            img = screenshot.convert("RGB").resize((50, 50))
            pixels = list(img.getdata())
            from collections import Counter
            most_common = Counter(pixels).most_common(5)
            lines = []
            for (r,g,b), cnt in most_common:
                lines.append(f"RGB({r},{g},{b}) = #{r:02X}{g:02X}{b:02X}  出現 {cnt} 次")
            return f"🎨 區域主要顏色：\n" + "\n".join(lines)
    except Exception as e:
        return f"❌ 顏色選取失敗：{e}"


def execute_webcam(action, duration=5.0, output="", device=0):
    try:
        import cv2, tempfile
        if action == "list":
            found = []
            for i in range(5):
                cap = cv2.VideoCapture(i)
                if cap.isOpened():
                    found.append(f"裝置 {i}")
                    cap.release()
            return f"📷 可用攝影機：\n" + ("\n".join(found) if found else "無")
        elif action == "photo":
            cap = cv2.VideoCapture(int(device))
            if not cap.isOpened():
                return f"❌ 無法開啟攝影機 {device}"
            ret, frame = cap.read()
            cap.release()
            if not ret:
                return "❌ 無法拍攝"
            out = output or str(Path("C:/Users/blue_/Desktop/測試檔案") / f"webcam_{datetime.datetime.now().strftime('%Y%m%d_%H%M%S')}.jpg")
            cv2.imwrite(out, frame)
            return f"✅ 已拍照：{out}"
        elif action == "video":
            cap = cv2.VideoCapture(int(device))
            if not cap.isOpened():
                return f"❌ 無法開啟攝影機 {device}"
            out = output or str(Path("C:/Users/blue_/Desktop/測試檔案") / f"webcam_{datetime.datetime.now().strftime('%Y%m%d_%H%M%S')}.avi")
            w = int(cap.get(cv2.CAP_PROP_FRAME_WIDTH))
            h = int(cap.get(cv2.CAP_PROP_FRAME_HEIGHT))
            writer = cv2.VideoWriter(out, cv2.VideoWriter_fourcc(*"XVID"), 20, (w,h))
            import time
            end = time.time() + float(duration)
            while time.time() < end:
                ret, frame = cap.read()
                if ret:
                    writer.write(frame)
            cap.release()
            writer.release()
            return f"✅ 已錄影 {duration}s：{out}"
    except Exception as e:
        return f"❌ 攝影機操作失敗：{e}"


def execute_multi_monitor(action, monitor=1, window=""):
    try:
        import subprocess, win32gui, win32con
        if action == "list":
            r = subprocess.run(["powershell", "-Command",
                "Get-CimInstance Win32_DesktopMonitor | Select-Object Name,ScreenWidth,ScreenHeight | Format-Table -AutoSize; (Get-CimInstance Win32_VideoController).CurrentHorizontalResolution"],
                capture_output=True, text=True, encoding="utf-8", errors="replace")
            return f"🖥️ 螢幕資訊：\n{r.stdout.strip()}"
        elif action in ("extend","clone"):
            mode = "/extend" if action == "extend" else "/clone"
            subprocess.run(["displayswitch.exe", mode])
            return f"✅ 螢幕模式已切換為：{action}"
        elif action == "set_primary":
            r = subprocess.run(["powershell", "-Command",
                f"$monitors = Get-CimInstance Win32_DesktopMonitor; Set-DisplayResolution -Force"],
                capture_output=True, text=True)
            return f"⚠️ 設定主螢幕需要更細部設定，目前螢幕數量已顯示"
        elif action == "move_window":
            import ctypes
            user32 = ctypes.windll.user32
            sw = user32.GetSystemMetrics(0)
            hwnds = []
            win32gui.EnumWindows(
                lambda h, l: l.append(h) if win32gui.IsWindowVisible(h) and window.lower() in win32gui.GetWindowText(h).lower() else None,
                hwnds)
            if not hwnds:
                return f"⚠️ 找不到視窗：{window}"
            offset_x = sw * (int(monitor) - 1)
            rect = win32gui.GetWindowRect(hwnds[0])
            win32gui.MoveWindow(hwnds[0], offset_x + 100, 100, rect[2]-rect[0], rect[3]-rect[1], True)
            return f"✅ 已移動視窗 '{window}' 到螢幕 {monitor}"
    except Exception as e:
        return f"❌ 多螢幕管理失敗：{e}"


def execute_printer(action, path="", printer_name=""):
    try:
        import subprocess, win32print
        if action == "list":
            printers = win32print.EnumPrinters(win32print.PRINTER_ENUM_LOCAL | win32print.PRINTER_ENUM_CONNECTIONS)
            lines = [f"- {p[2]}" for p in printers]
            default = win32print.GetDefaultPrinter()
            return f"🖨️ 印表機清單（預設：{default}）：\n" + "\n".join(lines)
        elif action == "print":
            pname = printer_name or win32print.GetDefaultPrinter()
            import win32api
            win32api.ShellExecute(0, "print", path, f'/d:"{pname}"', ".", 0)
            return f"✅ 已傳送列印：{path} → {pname}"
        elif action == "queue":
            pname = printer_name or win32print.GetDefaultPrinter()
            handle = win32print.OpenPrinter(pname)
            jobs = win32print.EnumJobs(handle, 0, -1, 1)
            win32print.ClosePrinter(handle)
            if not jobs:
                return f"✅ {pname} 列印佇列為空"
            lines = [f"工作 {j['JobId']}：{j['pDocument']} ({j['Status']})" for j in jobs]
            return "🖨️ 列印佇列：\n" + "\n".join(lines)
        elif action == "clear_queue":
            pname = printer_name or win32print.GetDefaultPrinter()
            handle = win32print.OpenPrinter(pname)
            jobs = win32print.EnumJobs(handle, 0, -1, 1)
            for j in jobs:
                win32print.SetJob(handle, j["JobId"], 0, None, win32print.JOB_CONTROL_DELETE)
            win32print.ClosePrinter(handle)
            return f"✅ 已清除 {pname} 的列印佇列"
        elif action == "set_default":
            win32print.SetDefaultPrinter(printer_name)
            return f"✅ 預設印表機已設為：{printer_name}"
    except Exception as e:
        return f"❌ 印表機操作失敗：{e}"


def execute_wifi(action, ssid="", password=""):
    try:
        import subprocess
        if action == "scan":
            r = subprocess.run(["netsh", "wlan", "show", "networks", "mode=Bssid"],
                capture_output=True, text=True, encoding="utf-8", errors="replace")
            return f"📡 附近 Wi-Fi：\n{r.stdout.strip()[:2000]}"
        elif action == "status":
            r = subprocess.run(["netsh", "wlan", "show", "interfaces"],
                capture_output=True, text=True, encoding="utf-8", errors="replace")
            return f"📡 Wi-Fi 狀態：\n{r.stdout.strip()}"
        elif action == "saved":
            r = subprocess.run(["netsh", "wlan", "show", "profiles"],
                capture_output=True, text=True, encoding="utf-8", errors="replace")
            return f"📡 已儲存 Wi-Fi：\n{r.stdout.strip()}"
        elif action == "password":
            r = subprocess.run(["netsh", "wlan", "show", "profile", f"name={ssid}", "key=clear"],
                capture_output=True, text=True, encoding="utf-8", errors="replace")
            return f"🔑 '{ssid}' 密碼資訊：\n{r.stdout.strip()}"
        elif action == "connect":
            r = subprocess.run(["netsh", "wlan", "connect", f"name={ssid}"],
                capture_output=True, text=True, encoding="utf-8", errors="replace")
            return f"✅ 嘗試連線 '{ssid}'：{r.stdout.strip()}"
        elif action == "disconnect":
            r = subprocess.run(["netsh", "wlan", "disconnect"],
                capture_output=True, text=True, encoding="utf-8", errors="replace")
            return f"✅ Wi-Fi 已斷線：{r.stdout.strip()}"
    except Exception as e:
        return f"❌ Wi-Fi 操作失敗：{e}"


def execute_proxy(action, host=""):
    try:
        import winreg
        key_path = r"Software\Microsoft\Windows\CurrentVersion\Internet Settings"
        with winreg.OpenKey(winreg.HKEY_CURRENT_USER, key_path, 0, winreg.KEY_READ | winreg.KEY_SET_VALUE) as k:
            if action == "get":
                try:
                    enabled = winreg.QueryValueEx(k, "ProxyEnable")[0]
                    server = winreg.QueryValueEx(k, "ProxyServer")[0]
                    return f"🌐 代理設定：{'啟用' if enabled else '停用'}\n伺服器：{server}"
                except Exception:
                    return "🌐 代理：未設定"
            elif action == "set":
                winreg.SetValueEx(k, "ProxyEnable", 0, winreg.REG_DWORD, 1)
                winreg.SetValueEx(k, "ProxyServer", 0, winreg.REG_SZ, host)
                return f"✅ 代理已設定：{host}"
            elif action == "disable":
                winreg.SetValueEx(k, "ProxyEnable", 0, winreg.REG_DWORD, 0)
                return "✅ 代理已停用"
    except Exception as e:
        return f"❌ 代理設定失敗：{e}"


def execute_lock_screen(action):
    try:
        import subprocess
        if action == "lock":
            subprocess.Popen(["rundll32.exe", "user32.dll,LockWorkStation"])
            return "🔒 螢幕已鎖定"
        elif action == "logoff":
            subprocess.run(["shutdown", "/l"], capture_output=True)
            return "✅ 已登出"
        elif action == "switch_user":
            subprocess.Popen(["tsdiscon.exe"])
            return "✅ 已切換使用者"
    except Exception as e:
        return f"❌ 鎖定/登出失敗：{e}"


def execute_defender(action, path=""):
    try:
        import subprocess
        if action == "status":
            r = subprocess.run(["powershell", "-Command",
                "Get-MpComputerStatus | Select-Object AMRunningMode,RealTimeProtectionEnabled,AntivirusSignatureLastUpdated | Format-List"],
                capture_output=True, text=True, encoding="utf-8", errors="replace")
            return f"🛡️ Defender 狀態：\n{r.stdout.strip()}"
        elif action == "quick_scan":
            r = subprocess.run(["powershell", "-Command",
                "Start-MpScan -ScanType QuickScan"],
                capture_output=True, text=True, timeout=30)
            return "🛡️ 快速掃描已啟動（背景執行中）"
        elif action == "full_scan":
            r = subprocess.run(["powershell", "-Command",
                "Start-MpScan -ScanType FullScan"],
                capture_output=True, text=True, timeout=30)
            return "🛡️ 完整掃描已啟動（背景執行中）"
        elif action == "threats":
            r = subprocess.run(["powershell", "-Command",
                "Get-MpThreatDetection | Select-Object ThreatID,Resources,ActionSuccess | Format-Table -AutoSize"],
                capture_output=True, text=True, encoding="utf-8", errors="replace")
            out = r.stdout.strip()
            return f"🛡️ 威脅記錄：\n{out}" if out else "✅ 無威脅記錄"
        elif action == "add_exclusion":
            r = subprocess.run(["powershell", "-Command",
                f"Add-MpPreference -ExclusionPath '{path}'"],
                capture_output=True, text=True)
            return f"✅ 已新增排除：{path}" if r.returncode == 0 else f"❌ 失敗：{r.stderr.strip()}"
        elif action == "remove_exclusion":
            r = subprocess.run(["powershell", "-Command",
                f"Remove-MpPreference -ExclusionPath '{path}'"],
                capture_output=True, text=True)
            return f"✅ 已移除排除：{path}" if r.returncode == 0 else f"❌ 失敗：{r.stderr.strip()}"
        elif action == "list_exclusions":
            r = subprocess.run(["powershell", "-Command",
                "(Get-MpPreference).ExclusionPath"],
                capture_output=True, text=True, encoding="utf-8", errors="replace")
            out = r.stdout.strip()
            return f"🛡️ 排除清單：\n{out}" if out else "✅ 無排除項目"
    except subprocess.TimeoutExpired:
        return "⏳ 掃描已啟動（在背景執行）"
    except Exception as e:
        return f"❌ Defender 操作失敗：{e}"


def execute_email_read(host, user, password, folder="INBOX", count=5):
    try:
        import imapclient, email as _email
        from email.header import decode_header
        client = imapclient.IMAPClient(host, ssl=True)
        client.login(user, password)
        client.select_folder(folder)
        msgs = client.search(["ALL"])
        recent = msgs[-count:] if len(msgs) >= count else msgs
        results = []
        for uid in reversed(recent):
            raw = client.fetch([uid], ["RFC822"])[uid][b"RFC822"]
            msg = _email.message_from_bytes(raw)
            subj_raw, enc = decode_header(msg["Subject"])[0]
            subject = subj_raw.decode(enc or "utf-8") if isinstance(subj_raw, bytes) else subj_raw
            body = ""
            if msg.is_multipart():
                for part in msg.walk():
                    if part.get_content_type() == "text/plain":
                        body = part.get_payload(decode=True).decode(errors="replace")[:150]; break
            else:
                body = msg.get_payload(decode=True).decode(errors="replace")[:150]
            results.append(f"寄件人：{msg['From']}\n主旨：{subject}\n日期：{msg['Date']}\n{body}\n{'─'*30}")
        client.logout()
        return "\n".join(results) if results else "（收件匣為空）"
    except Exception as e:
        return f"❌ 讀取郵件失敗：{e}"


def execute_calendar(action, days=7, title="", start="", end="", description=""):
    try:
        from google.oauth2.credentials import Credentials
        from googleapiclient.discovery import build
        from datetime import timezone, timedelta
        creds_path = Path("C:/Users/blue_/claude-telegram-bot/gcal_token.json")
        if not creds_path.exists():
            return "❌ 未找到 Google Calendar 憑證（gcal_token.json）"
        creds = Credentials.from_authorized_user_file(str(creds_path))
        service = build("calendar", "v3", credentials=creds)
        if action == "list":
            now = datetime.datetime.now(datetime.timezone.utc)
            events = service.events().list(
                calendarId="primary", timeMin=now.isoformat(),
                timeMax=(now + datetime.timedelta(days=days)).isoformat(),
                maxResults=20, singleEvents=True, orderBy="startTime"
            ).execute().get("items", [])
            if not events:
                return f"未來 {days} 天沒有行程"
            return "\n".join(f"📅 {e['start'].get('dateTime',e['start'].get('date'))}  {e.get('summary','（無標題）')}" for e in events)
        elif action == "add":
            event = {"summary": title, "description": description,
                     "start": {"dateTime": start, "timeZone": "Asia/Taipei"},
                     "end": {"dateTime": end, "timeZone": "Asia/Taipei"}}
            created = service.events().insert(calendarId="primary", body=event).execute()
            return f"✅ 行程已新增：{created.get('summary')}"
        elif action == "delete":
            if not title:
                return "❌ 請提供要刪除的行程標題"
            events = service.events().list(
                calendarId="primary", q=title,
                maxResults=5, singleEvents=True
            ).execute().get("items", [])
            if not events:
                return f"❌ 找不到符合的行程：{title}"
            service.events().delete(calendarId="primary", eventId=events[0]["id"]).execute()
            return f"✅ 已刪除行程：{events[0].get('summary', title)}"
    except Exception as e:
        return f"❌ 行事曆操作失敗：{e}"


def execute_global_hotkey(hotkey, command, duration=60.0):
    try:
        import keyboard as kb, time as t
        triggered = []
        def on_trigger():
            triggered.append(datetime.datetime.now().strftime("%H:%M:%S"))
            subprocess.run(command, shell=True)
        kb.add_hotkey(hotkey, on_trigger)
        t.sleep(duration)
        kb.remove_all_hotkeys()
        return f"✅ 快捷鍵 [{hotkey}] 共觸發 {len(triggered)} 次"
    except Exception as e:
        return f"❌ 快捷鍵監聽失敗：{e}"


def execute_git(action, repo=".", message="", branch="master"):
    try:
        import git as _git
        r = _git.Repo(repo)
        if action == "status": return r.git.status()
        elif action == "log":
            return "\n".join(f"{c.hexsha[:7]} [{c.authored_datetime.strftime('%m-%d %H:%M')}] {c.message.strip()[:60]}" for c in list(r.iter_commits())[:10])
        elif action == "pull":
            result = r.remotes.origin.pull(); return f"✅ Pull 完成"
        elif action == "add": r.git.add(A=True); return "✅ git add -A 完成"
        elif action == "commit": r.index.commit(message or "auto commit"); return f"✅ committed: {message}"
        elif action == "push": r.remotes.origin.push(branch); return f"✅ pushed to origin/{branch}"
        elif action == "diff": return r.git.diff()[:2000] or "（無變更）"
        return "未知動作"
    except Exception as e:
        return f"❌ Git 失敗：{e}"


def execute_hardware():
    try:
        import psutil
        battery = psutil.sensors_battery()
        bat = f"{battery.percent:.0f}% {'充電中' if battery.power_plugged else '使用電池'}" if battery else "無電池"
        gpu_str = ""
        try:
            import GPUtil
            gpus = GPUtil.getGPUs()
            gpu_str = "\n".join(f"GPU {g.name}: 使用率 {g.load*100:.0f}% | 記憶體 {g.memoryUsed:.0f}/{g.memoryTotal:.0f}MB | 溫度 {g.temperature}°C" for g in gpus)
        except Exception:
            gpu_str = "GPU：未偵測到 NVIDIA GPU"
        return f"🔋 電池：{bat}\n{gpu_str}"
    except Exception as e:
        return f"❌ 硬體監控失敗：{e}"


def execute_report(title, data_json, output=""):
    try:
        import jinja2, json
        data = json.loads(data_json)
        out_path = output or str(Path("C:/Users/blue_/Desktop/測試檔案") / f"report_{datetime.datetime.now().strftime('%Y%m%d_%H%M%S')}.html")
        tmpl = jinja2.Template("""<!DOCTYPE html><html><head><meta charset="utf-8">
<style>body{font-family:sans-serif;margin:40px}table{border-collapse:collapse;width:100%}
th,td{border:1px solid #ccc;padding:8px}th{background:#4472C4;color:white}
tr:nth-child(even){background:#f2f2f2}h1{color:#4472C4}</style></head>
<body><h1>{{ title }}</h1><p>生成時間：{{ time }}</p>
{% for section, rows in data.items() %}<h2>{{ section }}</h2>
{% if rows is iterable and rows is not string %}{% if rows[0] is mapping %}
<table><tr>{% for k in rows[0].keys() %}<th>{{ k }}</th>{% endfor %}</tr>
{% for row in rows %}<tr>{% for v in row.values() %}<td>{{ v }}</td>{% endfor %}</tr>{% endfor %}</table>
{% else %}<ul>{% for i in rows %}<li>{{ i }}</li>{% endfor %}</ul>{% endif %}
{% else %}<p>{{ rows }}</p>{% endif %}{% endfor %}</body></html>""")
        Path(out_path).write_text(tmpl.render(title=title, data=data, time=datetime.datetime.now().strftime("%Y-%m-%d %H:%M:%S")), encoding="utf-8")
        return f"✅ 報告已生成：{out_path}"
    except Exception as e:
        return f"❌ 報告生成失敗：{e}"


def execute_dropbox(action, local, remote, token=""):
    try:
        import dropbox as dbx
        tok = token or os.getenv("DROPBOX_TOKEN", "")
        if not tok: return "❌ 請設定 DROPBOX_TOKEN 環境變數"
        d = dbx.Dropbox(tok)
        if action == "upload":
            with open(local, "rb") as f:
                d.files_upload(f.read(), remote, mode=dbx.files.WriteMode.overwrite)
            return f"✅ 已上傳到 Dropbox：{remote}"
        elif action == "download":
            _, res = d.files_download(remote)
            Path(local).write_bytes(res.content)
            return f"✅ 已從 Dropbox 下載：{local}"
    except Exception as e:
        return f"❌ Dropbox 失敗：{e}"


def execute_docker(action, name=""):
    try:
        import docker as _docker
        client = _docker.from_env()
        if action == "list":
            return "\n".join(f"[{c.status}] {c.name} {c.image.tags}" for c in client.containers.list(all=True)) or "（無容器）"
        elif action == "start": client.containers.get(name).start(); return f"✅ {name} 已啟動"
        elif action == "stop": client.containers.get(name).stop(); return f"✅ {name} 已停止"
        elif action == "logs": return client.containers.get(name).logs(tail=50).decode(errors="replace")
        elif action == "images": return "\n".join(f"{img.tags} {img.short_id}" for img in client.images.list())
    except Exception as e:
        return f"❌ Docker 失敗：{e}"


def execute_pdf_to_image(path, output_dir="", dpi=150):
    try:
        import fitz
        doc = fitz.open(path)
        out = Path(output_dir) if output_dir else Path(path).parent / (Path(path).stem + "_imgs")
        out.mkdir(parents=True, exist_ok=True)
        for i, page in enumerate(doc):
            pix = page.get_pixmap(matrix=fitz.Matrix(dpi/72, dpi/72))
            pix.save(str(out / f"page_{i+1}.png"))
        return f"✅ 已轉換 {len(doc)} 頁到：{out}"
    except Exception as e:
        return f"❌ PDF 轉圖片失敗：{e}"


def execute_barcode(image_path=""):
    try:
        from pyzbar.pyzbar import decode
        from PIL import Image
        img = Image.open(image_path) if image_path else pyautogui.screenshot()
        results = decode(img)
        if not results: return "❌ 未偵測到條碼"
        return "\n".join(f"類型：{r.type}  內容：{r.data.decode('utf-8', errors='replace')}" for r in results)
    except Exception as e:
        return f"❌ 條碼掃描失敗：{e}"


def execute_nlp(action, text):
    try:
        import anthropic as _ant
        c = _ant.Anthropic(api_key=os.getenv("ANTHROPIC_API_KEY"))
        prompt = f"請用繁體中文摘要以下文字（100字以內）：\n\n{text}" if action == "summarize" else f"分析以下文字的情緒，只回覆：正面/負面/中性 + 一句說明：\n\n{text}"
        msg = c.messages.create(model="claude-haiku-4-5-20251001", max_tokens=256, messages=[{"role":"user","content":prompt}])
        return msg.content[0].text
    except Exception as e:
        return f"❌ NLP 失敗：{e}"


def execute_vpn(action, name="", user="", password=""):
    try:
        if action == "list":
            r = subprocess.run(["powershell.exe","-Command","Get-VpnConnection | Select-Object Name,ConnectionStatus | Format-Table"], capture_output=True, text=True, encoding="utf-8", errors="replace")
            return r.stdout or "（未設定 VPN）"
        elif action == "connect":
            r = subprocess.run(["rasdial", name, user, password], capture_output=True, text=True, encoding="cp950", errors="replace")
            return r.stdout.strip()
        elif action == "disconnect":
            r = subprocess.run(["rasdial", name, "/disconnect"], capture_output=True, text=True, encoding="cp950", errors="replace")
            return r.stdout.strip()
    except Exception as e:
        return f"❌ VPN 失敗：{e}"


def execute_restore_point(action, description=""):
    try:
        if action == "create":
            ps = f"Checkpoint-Computer -Description '{description or 'Claude Auto Restore'}' -RestorePointType MODIFY_SETTINGS"
            r = subprocess.run(["powershell.exe","-Command",ps], capture_output=True, text=True, encoding="utf-8", errors="replace")
            return r.stdout or "✅ 還原點已建立"
        elif action == "list":
            r = subprocess.run(["powershell.exe","-Command","Get-ComputerRestorePoint | Select-Object SequenceNumber,Description,CreationTime | Format-Table"], capture_output=True, text=True, encoding="utf-8", errors="replace")
            return r.stdout or "（無還原點）"
    except Exception as e:
        return f"❌ 系統還原點失敗：{e}"


def execute_disk_analyze(path="C:/", top=10):
    try:
        import psutil
        usage = psutil.disk_usage(path)
        result = f"磁碟：{path}\n總容量：{usage.total/1024**3:.1f} GB | 已使用：{usage.used/1024**3:.1f} GB ({usage.percent}%) | 可用：{usage.free/1024**3:.1f} GB\n\n"
        sizes = []
        for item in Path(path).iterdir():
            try:
                size = sum(f.stat().st_size for f in item.rglob("*") if f.is_file()) if item.is_dir() else item.stat().st_size
                sizes.append((size, str(item)))
            except Exception: pass
        sizes.sort(reverse=True)
        result += "\n".join(f"{s/1024**3:.2f} GB  {n}" for s, n in sizes[:top])
        return result
    except Exception as e:
        return f"❌ 磁碟分析失敗：{e}"


def execute_face_detect(image_path="", output=""):
    try:
        import cv2, numpy as np
        img = cv2.cvtColor(np.array(pyautogui.screenshot()), cv2.COLOR_RGB2BGR) if not image_path else cv2.imread(image_path)
        cascade = cv2.CascadeClassifier(cv2.data.haarcascades + "haarcascade_frontalface_default.xml")
        faces = cascade.detectMultiScale(cv2.cvtColor(img, cv2.COLOR_BGR2GRAY), 1.1, 4)
        for (x, y, w, h) in faces:
            cv2.rectangle(img, (x, y), (x+w, y+h), (0, 255, 0), 2)
        out_path = output or str(Path("C:/Users/blue_/Desktop/測試檔案") / f"faces_{datetime.datetime.now().strftime('%H%M%S')}.jpg")
        cv2.imwrite(out_path, img)
        return f"✅ 偵測到 {len(faces)} 張人臉：{out_path}"
    except Exception as e:
        return f"❌ 人臉偵測失敗：{e}"


def execute_video_gif(path, start=0, duration=5.0, output="", fps=10):
    try:
        import imageio
        out = output or path.replace(".mp4", ".gif")
        reader = imageio.get_reader(path)
        video_fps = reader.get_meta_data().get("fps", 30)
        frames = [f for i, f in enumerate(reader) if int(start*video_fps) <= i < int((start+duration)*video_fps)]
        imageio.mimsave(out, frames, fps=fps)
        return f"✅ GIF 已生成：{out}（{len(frames)} 幀）"
    except Exception as e:
        return f"❌ 影片轉 GIF 失敗：{e}"


def execute_excel_chart(path, sheet, chart_type="bar", title=""):
    try:
        import openpyxl
        from openpyxl.chart import BarChart, LineChart, PieChart, Reference
        wb = openpyxl.load_workbook(path)
        ws = wb[sheet] if sheet in wb.sheetnames else wb.active
        chart = {"bar": BarChart, "line": LineChart, "pie": PieChart}.get(chart_type, BarChart)()
        chart.title = title or sheet; chart.style = 10
        data = Reference(ws, min_col=2, min_row=1, max_row=ws.max_row, max_col=ws.max_column)
        chart.add_data(data, titles_from_data=True)
        chart.set_categories(Reference(ws, min_col=1, min_row=2, max_row=ws.max_row))
        ws.add_chart(chart, "A" + str(ws.max_row + 2))
        wb.save(path); return f"✅ 圖表已加入：{path}"
    except Exception as e:
        return f"❌ Excel 圖表失敗：{e}"


def execute_speedtest():
    try:
        import speedtest as st
        s = st.Speedtest(); s.get_best_server()
        dl = s.download()/1_000_000; ul = s.upload()/1_000_000
        return f"📶 下載：{dl:.1f} Mbps | 上傳：{ul:.1f} Mbps | Ping：{s.results.ping:.0f} ms | 伺服器：{s.results.server.get('name','')}"
    except Exception as e:
        return f"❌ 速度測試失敗：{e}"


def execute_screenshot_compare(img1_path="", img2_path="", output=""):
    try:
        import cv2, numpy as np, time as t
        img1 = cv2.cvtColor(np.array(pyautogui.screenshot()), cv2.COLOR_RGB2BGR) if not img1_path else cv2.imread(img1_path)
        if not img2_path: t.sleep(2); img2 = cv2.cvtColor(np.array(pyautogui.screenshot()), cv2.COLOR_RGB2BGR)
        else: img2 = cv2.imread(img2_path)
        h = min(img1.shape[0], img2.shape[0]); w = min(img1.shape[1], img2.shape[1])
        diff = cv2.absdiff(img1[:h,:w], img2[:h,:w])
        _, thresh = cv2.threshold(cv2.cvtColor(diff, cv2.COLOR_BGR2GRAY), 30, 255, cv2.THRESH_BINARY)
        contours, _ = cv2.findContours(thresh, cv2.RETR_EXTERNAL, cv2.CHAIN_APPROX_SIMPLE)
        result = img2[:h,:w].copy(); cv2.drawContours(result, contours, -1, (0, 0, 255), 2)
        out = output or str(Path("C:/Users/blue_/Desktop/測試檔案") / f"diff_{datetime.datetime.now().strftime('%H%M%S')}.png")
        cv2.imwrite(out, result)
        pct = cv2.countNonZero(thresh) / (h * w) * 100
        return f"差異：{pct:.2f}%，標記圖：{out}"
    except Exception as e:
        return f"❌ 截圖比對失敗：{e}"


def execute_screen_record(action, duration=10.0, output=""):
    try:
        if action == "record":
            import mss, cv2, numpy as np, time as t
            out_path = output or str(Path("C:/Users/blue_/Desktop/測試檔案") / f"record_{datetime.datetime.now().strftime('%H%M%S')}.mp4")
            with mss.mss() as sct:
                mon = sct.monitors[1]
                w, h = mon["width"], mon["height"]
                writer = cv2.VideoWriter(out_path, cv2.VideoWriter_fourcc(*"mp4v"), 10, (w, h))
                end = t.time() + duration
                while t.time() < end:
                    frame = np.array(sct.grab(mon))
                    writer.write(cv2.cvtColor(frame, cv2.COLOR_BGRA2BGR))
                    t.sleep(0.1)
                writer.release()
            return f"✅ 錄影完成：{out_path}"
        elif action == "webcam":
            import cv2
            out_path = output or str(Path("C:/Users/blue_/Desktop/測試檔案") / f"webcam_{datetime.datetime.now().strftime('%H%M%S')}.jpg")
            cap = cv2.VideoCapture(0)
            ret, frame = cap.read()
            cap.release()
            if ret:
                cv2.imwrite(out_path, frame)
                return f"✅ 已拍照：{out_path}"
            return "❌ 無法存取攝影機"
    except Exception as e:
        return f"❌ 失敗：{e}"


def execute_translate(text, target="zh-TW", source="auto"):
    try:
        from deep_translator import GoogleTranslator
        return GoogleTranslator(source=source, target=target).translate(text)
    except Exception as e:
        return f"❌ 翻譯失敗：{e}"


def execute_chart(chart_type, data_json, title="", output=""):
    try:
        import matplotlib
        matplotlib.use("Agg")
        import matplotlib.pyplot as plt, json
        data = json.loads(data_json)
        out_path = output or str(Path("C:/Users/blue_/Desktop/測試檔案") / f"chart_{datetime.datetime.now().strftime('%H%M%S')}.png")
        fig, ax = plt.subplots()
        if chart_type == "line":
            for label, values in data.items():
                ax.plot(values, label=label)
            ax.legend()
        elif chart_type == "bar":
            ax.bar(list(data.keys()), list(data.values()))
        elif chart_type == "pie":
            ax.pie(list(data.values()), labels=list(data.keys()), autopct="%1.1f%%")
        if title:
            ax.set_title(title)
        plt.tight_layout()
        plt.savefig(out_path)
        plt.close()
        return out_path
    except Exception as e:
        return f"❌ 圖表生成失敗：{e}"


def execute_pptx(action, path, slides=""):
    try:
        from pptx import Presentation
        if action == "read":
            prs = Presentation(path)
            lines = []
            for i, slide in enumerate(prs.slides, 1):
                texts = [sh.text for sh in slide.shapes if sh.has_text_frame]
                lines.append(f"[投影片 {i}] " + " | ".join(t for t in texts if t.strip()))
            return "\n".join(lines) or "（簡報為空）"
        elif action == "create":
            import json
            from pptx.util import Pt
            data = json.loads(slides)
            prs = Presentation()
            for s in data:
                sl = prs.slides.add_slide(prs.slide_layouts[1])
                if "title" in s:
                    sl.shapes.title.text = s["title"]
                if "body" in s:
                    sl.placeholders[1].text = s["body"]
            prs.save(path)
            return f"✅ 已建立簡報：{path}"
    except Exception as e:
        return f"❌ PPT 操作失敗：{e}"


def execute_api_call(method, url, headers="{}", body="{}"):
    try:
        import json
        h = json.loads(headers) if headers else {}
        b = json.loads(body) if body else None
        resp = requests.request(method.upper(), url, headers=h, json=b, timeout=30)
        try:
            return json.dumps(resp.json(), ensure_ascii=False, indent=2)[:2000]
        except Exception:
            return resp.text[:2000]
    except Exception as e:
        return f"❌ API 呼叫失敗：{e}"


def execute_watchdog(process, script, duration=60.0):
    import psutil, time as t
    restarts = 0
    end = t.time() + duration
    while t.time() < end:
        running = any(p.name().lower() == process.lower() for p in psutil.process_iter())
        if not running:
            subprocess.Popen(["pythonw", script] if script.endswith(".py") else [script])
            restarts += 1
        t.sleep(5)
    return f"守護結束，共重啟 {restarts} 次"


def execute_ssh_sftp(action, host, user, password, command="", local="", remote="", port=22):
    try:
        import paramiko
        if action == "ssh_run":
            c = paramiko.SSHClient()
            c.set_missing_host_key_policy(paramiko.AutoAddPolicy())
            c.connect(host, port=port, username=user, password=password, timeout=15)
            _, stdout, stderr = c.exec_command(command)
            out = stdout.read().decode(errors="replace") + stderr.read().decode(errors="replace")
            c.close()
            return out.strip() or "（執行完畢，無輸出）"
        else:
            t = paramiko.Transport((host, port))
            t.connect(username=user, password=password)
            sftp = paramiko.SFTPClient.from_transport(t)
            if action == "sftp_upload":
                sftp.put(local, remote)
                result = f"✅ 已上傳：{local} → {remote}"
            else:
                sftp.get(remote, local)
                result = f"✅ 已下載：{remote} → {local}"
            sftp.close(); t.close()
            return result
    except Exception as e:
        return f"❌ SSH/SFTP 失敗：{e}"


def execute_network_diag(action, host, ports="22,80,443,3306,3389,8080"):
    try:
        if action == "ping":
            r = subprocess.run(["ping", "-n", "4", host], capture_output=True, text=True, encoding="cp950", errors="replace", timeout=20)
            return r.stdout.strip()
        elif action == "traceroute":
            r = subprocess.run(["tracert", host], capture_output=True, text=True, encoding="cp950", errors="replace", timeout=60)
            return r.stdout[:2000]
        elif action == "portscan":
            import socket
            results = []
            for p in [int(x) for x in ports.split(",")]:
                s = socket.socket(); s.settimeout(1)
                r = s.connect_ex((host, p))
                results.append(f"Port {p}: {'開放 ✅' if r == 0 else '關閉 ❌'}")
                s.close()
            return "\n".join(results)
    except Exception as e:
        return f"❌ 網路診斷失敗：{e}"


def execute_win_service(action, name=""):
    try:
        if action == "list":
            r = subprocess.run(["powershell.exe", "-Command",
                "Get-Service | Select-Object Name,Status | Format-Table -AutoSize"],
                capture_output=True, text=True, encoding="utf-8", errors="replace")
            return r.stdout[:2000]
        else:
            cmd = f"{'Start' if action=='start' else 'Stop'}-Service -Name '{name}' -Force"
            r = subprocess.run(["powershell.exe", "-Command", cmd],
                capture_output=True, text=True, encoding="utf-8", errors="replace")
            return r.stdout or r.stderr or f"✅ {action} {name}"
    except Exception as e:
        return f"❌ 服務操作失敗：{e}"


def execute_pdf_edit(action, path="", output="", paths="", text=""):
    try:
        import fitz, json
        if action == "merge":
            pdf_list = json.loads(paths)
            writer = fitz.open()
            for p in pdf_list:
                writer.insert_pdf(fitz.open(p))
            writer.save(output)
            return f"✅ 已合併 {len(pdf_list)} 個 PDF：{output}"
        elif action == "split":
            doc = fitz.open(path)
            out_dir = Path(output)
            out_dir.mkdir(parents=True, exist_ok=True)
            for i, _ in enumerate(doc):
                out = fitz.open()
                out.insert_pdf(doc, from_page=i, to_page=i)
                out.save(str(out_dir / f"page_{i+1}.pdf"))
            return f"✅ 已分割 {len(doc)} 頁到：{output}"
        elif action == "watermark":
            doc = fitz.open(path)
            out_path = output or path.replace(".pdf", "_wm.pdf")
            for page in doc:
                page.insert_text((page.rect.width/2-50, page.rect.height/2),
                    text, fontsize=40, color=(0.8,0.8,0.8), rotate=45)
            doc.save(out_path)
            return f"✅ 已加浮水印：{out_path}"
    except Exception as e:
        return f"❌ PDF 編輯失敗：{e}"


def execute_audio_process(action, input_path, output="", start_ms=0, end_ms=0):
    try:
        from pydub import AudioSegment
        if action == "convert":
            fmt = Path(output).suffix.lstrip(".")
            AudioSegment.from_file(input_path).export(output, format=fmt)
            return f"✅ 已轉換：{output}"
        elif action == "trim":
            audio = AudioSegment.from_file(input_path)[start_ms:end_ms]
            out = output or input_path.replace(".", "_trim.")
            audio.export(out, format=Path(out).suffix.lstrip("."))
            return f"✅ 已剪輯：{out}"
    except Exception as e:
        return f"❌ 音訊處理失敗：{e}"


def execute_push_notify(platform, message, webhook_or_token):
    try:
        if platform == "discord":
            resp = requests.post(webhook_or_token, json={"content": message}, timeout=10)
            return f"✅ Discord 已發送（{resp.status_code}）"
        elif platform == "line":
            resp = requests.post(
                "https://notify-api.line.me/api/notify",
                headers={"Authorization": f"Bearer {webhook_or_token}"},
                data={"message": message}, timeout=10
            )
            return f"✅ LINE 已發送（{resp.status_code}）"
    except Exception as e:
        return f"❌ 推播失敗：{e}"


def execute_disk_backup(action, src="", dest=""):
    try:
        import tempfile, shutil
        tmp = Path(tempfile.gettempdir())
        if action == "list_temp":
            files = list(tmp.rglob("*"))
            total = sum(f.stat().st_size for f in files if f.is_file())
            return f"暫存資料夾：{tmp}\n檔案數：{len(files)}\n佔用空間：{total/1024/1024:.1f} MB"
        elif action == "clean_temp":
            count = 0
            for f in tmp.iterdir():
                try:
                    if f.is_file(): f.unlink(); count += 1
                    elif f.is_dir(): shutil.rmtree(f, ignore_errors=True); count += 1
                except Exception: pass
            return f"✅ 已清理 {count} 個暫存項目"
        elif action == "backup":
            out = Path(dest) / f"{Path(src).name}_{datetime.datetime.now().strftime('%Y%m%d_%H%M%S')}"
            shutil.make_archive(str(out), "zip", src)
            return f"✅ 備份完成：{out}.zip"
    except Exception as e:
        return f"❌ 磁碟操作失敗：{e}"


def execute_registry(action, key, value_name="", value=""):
    try:
        import winreg
        parts = key.split("\\", 1)
        roots = {"HKEY_LOCAL_MACHINE": winreg.HKEY_LOCAL_MACHINE,
                 "HKEY_CURRENT_USER": winreg.HKEY_CURRENT_USER,
                 "HKLM": winreg.HKEY_LOCAL_MACHINE, "HKCU": winreg.HKEY_CURRENT_USER}
        root = roots[parts[0]]
        if action == "read":
            with winreg.OpenKey(root, parts[1]) as k:
                if value_name:
                    val, _ = winreg.QueryValueEx(k, value_name)
                    return f"{value_name} = {val}"
                lines = []
                i = 0
                while True:
                    try:
                        n, v, _ = winreg.EnumValue(k, i); lines.append(f"{n} = {v}"); i += 1
                    except OSError: break
                return "\n".join(lines[:20])
        elif action == "write":
            with winreg.OpenKey(root, parts[1], 0, winreg.KEY_SET_VALUE) as k:
                winreg.SetValueEx(k, value_name, 0, winreg.REG_SZ, value)
            return f"✅ 已寫入：{value_name} = {value}"
    except Exception as e:
        return f"❌ 登錄檔操作失敗：{e}"


def execute_video_process(action, path, second=0, start=0, end=0, output=""):
    try:
        import cv2
        if action == "screenshot":
            cap = cv2.VideoCapture(path)
            cap.set(cv2.CAP_PROP_POS_MSEC, second * 1000)
            ret, frame = cap.read()
            cap.release()
            out = output or path.replace(".mp4", f"_frame{int(second)}s.jpg")
            if ret:
                cv2.imwrite(out, frame)
                return f"✅ 已擷取畫面：{out}"
            return "❌ 無法讀取影片"
        elif action == "trim":
            out = output or path.replace(".mp4", "_trim.mp4")
            subprocess.run(["ffmpeg", "-y", "-i", path, "-ss", str(start), "-to", str(end), "-c", "copy", out], capture_output=True)
            return f"✅ 已剪輯：{out}"
    except Exception as e:
        return f"❌ 影片處理失敗：{e}"


def execute_ai_video(prompt: str, provider: str = "replicate",
                     model: str = "", image_url: str = "",
                     duration: float = 5, output: str = "") -> str:
    """
    用 AI API 生成影片
    provider: replicate / runway / kling
    需要 .env 中對應的 API Key
    """
    import requests, time, traceback
    from pathlib import Path

    out = output or str(Path("C:/Users/blue_/Desktop/測試檔案") / f"ai_video_{int(time.time())}.mp4")

    def _download(url: str, dest: str) -> str:
        r = requests.get(url, timeout=120, stream=True)
        r.raise_for_status()
        with open(dest, "wb") as f:
            for chunk in r.iter_content(8192):
                f.write(chunk)
        return dest

    try:
        # ── Replicate ──────────────────────────────────────────────
        if provider == "replicate":
            api_key = os.getenv("REPLICATE_API_TOKEN", "")
            if not api_key:
                return "❌ 缺少 REPLICATE_API_TOKEN，請在 .env 加入"

            # 預設模型：minimax/video-01（Hailuo，高品質文字轉影片）
            mdl = model or ("stability-ai/stable-video-diffusion" if image_url else "minimax/video-01")

            headers = {
                "Authorization": f"Bearer {api_key}",
                "Content-Type": "application/json"
            }

            payload: dict = {"version": mdl} if "/" not in mdl else {}
            # 用 model 端點
            inputs: dict = {"prompt": prompt}
            if image_url:
                inputs["image"] = image_url
            if duration:
                inputs["duration"] = int(duration)

            # Replicate v1 API
            r = requests.post(
                f"https://api.replicate.com/v1/models/{mdl}/predictions" if "/" in mdl
                else f"https://api.replicate.com/v1/predictions",
                json={"input": inputs, **({"version": mdl} if "/" not in mdl else {})},
                headers=headers, timeout=120
            )
            r.raise_for_status()
            pred = r.json()
            if "id" not in pred:
                return f"❌ Replicate 回傳異常：{pred}"
            pred_id = pred["id"]
            pred_url = f"https://api.replicate.com/v1/predictions/{pred_id}"

            # 輪詢結果（最多等 10 分鐘，每 10 秒一次）
            for _ in range(60):
                time.sleep(10)
                resp = requests.get(pred_url, headers=headers, timeout=15)
                resp.raise_for_status()
                data = resp.json()
                status = data.get("status")
                if status == "succeeded":
                    video_url = data.get("output")
                    if isinstance(video_url, list):
                        video_url = video_url[0]
                    _download(video_url, out)
                    return f"✅ Replicate 影片已生成：{out}"
                elif status == "failed":
                    err = data.get("error", "未知錯誤")
                    return f"❌ Replicate 生成失敗：{err}"

            return "❌ Replicate 逾時（超過 10 分鐘）"

        # ── Runway ─────────────────────────────────────────────────
        elif provider == "runway":
            api_key = os.getenv("RUNWAY_API_KEY", "")
            if not api_key:
                return "❌ 缺少 RUNWAY_API_KEY，請在 .env 加入"
            try:
                import runwayml
            except ImportError:
                return "❌ 請先安裝：pip install runwayml"

            client = runwayml.RunwayML(api_key=api_key)

            if image_url:
                task = client.image_to_video.create(
                    model="gen4_turbo",
                    prompt_image=image_url,
                    prompt_text=prompt,
                    duration=int(min(duration, 10)),
                    ratio="1280:720"
                )
            else:
                # Runway Gen-4 需要圖片，改用 gen3a_turbo text-to-video
                task = client.text_to_video.create(
                    model="gen4_turbo",
                    prompt_text=prompt,
                    duration=int(min(duration, 10)),
                    ratio="1280:720"
                )

            task_id = task.id
            # 輪詢
            for _ in range(60):
                time.sleep(5)
                t = client.tasks.retrieve(task_id)
                if t.status == "SUCCEEDED":
                    _download(t.output[0], out)
                    return f"✅ Runway 影片已生成：{out}"
                elif t.status in ("FAILED", "CANCELLED"):
                    return f"❌ Runway 生成失敗：{t.failure_reason}"
            return "❌ Runway 逾時"

        # ── Kling（快影/可靈）────────────────────────────────────────
        elif provider == "kling":
            import hmac, hashlib, base64, json as _json
            access_key = os.getenv("KLING_ACCESS_KEY", "")
            secret_key = os.getenv("KLING_SECRET_KEY", "")
            if not access_key or not secret_key:
                return "❌ 缺少 KLING_ACCESS_KEY / KLING_SECRET_KEY，請在 .env 加入"

            # 生成 JWT
            import jwt as _jwt
            payload_jwt = {
                "iss": access_key,
                "exp": int(time.time()) + 1800,
                "nbf": int(time.time()) - 5
            }
            token = _jwt.encode(payload_jwt, secret_key, algorithm="HS256")
            headers = {"Authorization": f"Bearer {token}", "Content-Type": "application/json"}

            # 建立任務
            body: dict = {
                "model_name": "kling-v1",
                "prompt": prompt,
                "duration": str(int(min(duration, 10))),
                "aspect_ratio": "16:9"
            }
            if image_url:
                body["image_url"] = image_url

            endpoint = "image2video" if image_url else "text2video"
            r = requests.post(
                f"https://api.klingai.com/v1/videos/{endpoint}",
                json=body, headers=headers, timeout=30
            )
            r.raise_for_status()
            data = r.json()
            if data.get("code") != 0:
                return f"❌ Kling API 錯誤：{data.get('message')}"

            task_id = data["data"]["task_id"]

            # 輪詢
            for _ in range(60):
                time.sleep(5)
                # 重新生成 token
                payload_jwt["exp"] = int(time.time()) + 1800
                token = _jwt.encode(payload_jwt, secret_key, algorithm="HS256")
                headers["Authorization"] = f"Bearer {token}"

                resp = requests.get(
                    f"https://api.klingai.com/v1/videos/{endpoint}/{task_id}",
                    headers=headers, timeout=15
                )
                resp.raise_for_status()
                d = resp.json().get("data", {})
                status = d.get("task_status")
                if status == "succeed":
                    video_url = d["task_result"]["videos"][0]["url"]
                    _download(video_url, out)
                    return f"✅ Kling 影片已生成：{out}"
                elif status == "failed":
                    return f"❌ Kling 生成失敗：{d.get('task_status_msg')}"

            return "❌ Kling 逾時"

        else:
            return f"❌ 未知 provider：{provider}，可用：replicate / runway / kling"

    except Exception as e:
        return f"❌ AI 影片生成失敗：{e}\n{traceback.format_exc()}"


def execute_video_gen(mode: str = "slideshow", output: str = "", **kwargs) -> str:
    """
    生成影片
    mode: slideshow / text_video / tts_video / screen_record
    """
    import re as _re
    import numpy as np
    import subprocess
    import tempfile
    import traceback
    from pathlib import Path
    from PIL import Image, ImageDraw, ImageFont
    import imageio_ffmpeg

    ffmpeg_exe = imageio_ffmpeg.get_ffmpeg_exe()
    w, h = kwargs.get("size", (1280, 720))
    fps  = kwargs.get("fps", 24)
    out  = output or str(Path("C:/Users/blue_/Desktop/測試檔案") / f"video_{datetime.datetime.now().strftime('%H%M%S')}.mp4")

    def _write_frames(frames_iter, out_path, vid_fps, width, height):
        """ffmpeg pipe 寫 mp4，失敗時 raise"""
        proc = subprocess.Popen([
            ffmpeg_exe, "-y",
            "-f", "rawvideo", "-vcodec", "rawvideo",
            "-s", f"{width}x{height}", "-pix_fmt", "rgb24",
            "-r", str(vid_fps), "-i", "pipe:0",
            "-vcodec", "libx264", "-pix_fmt", "yuv420p",
            "-preset", "fast", "-crf", "23",
            out_path
        ], stdin=subprocess.PIPE, stderr=subprocess.PIPE)
        for frame in frames_iter:
            proc.stdin.write(np.asarray(frame, dtype=np.uint8).tobytes())
        proc.stdin.close()
        rc  = proc.wait()
        err = proc.stderr.read().decode(errors="replace")
        if rc != 0 or not Path(out_path).exists():
            raise RuntimeError(f"ffmpeg 失敗 rc={rc}: {err[-300:]}")

    def _get_font(pt=60):
        for fp in ["C:/Windows/Fonts/msjh.ttc", "C:/Windows/Fonts/msyh.ttc",
                   "C:/Windows/Fonts/simhei.ttf", "C:/Windows/Fonts/arial.ttf"]:
            if Path(fp).exists():
                try:
                    return ImageFont.truetype(fp, pt)
                except Exception:
                    pass
        return ImageFont.load_default()

    def _tts_sync(text: str, voice: str, out_path: str):
        """在當前 thread 建立新的 event loop 執行 edge-tts（thread 內無 running loop，安全）"""
        import asyncio
        import edge_tts
        loop = asyncio.new_event_loop()
        asyncio.set_event_loop(loop)
        try:
            comm = edge_tts.Communicate(text, voice, rate="-5%", pitch="-5Hz")
            loop.run_until_complete(comm.save(out_path))
        finally:
            loop.close()
        if not Path(out_path).exists():
            raise RuntimeError("TTS 語音檔案未生成")

    def _get_audio_dur(audio_path: str) -> float:
        r = subprocess.run([ffmpeg_exe, "-i", audio_path],
                           capture_output=True, text=True, errors="replace")
        m = _re.search(r"Duration:\s*(\d+):(\d+):(\d+\.?\d*)", r.stderr)
        if m:
            return int(m.group(1)) * 3600 + int(m.group(2)) * 60 + float(m.group(3))
        return 10.0

    try:
        # ── 投影片 ────────────────────────────────────────────────────
        if mode == "slideshow":
            images   = kwargs.get("images", [])
            dur      = kwargs.get("duration", 3)
            sl_fps   = kwargs.get("fps", 12)
            trans    = kwargs.get("transition", 0.5)
            if not images:
                return "❌ 請提供 images 參數（圖片路徑列表）"

            sf = int(sl_fps * dur)
            tf = int(sl_fps * trans)

            def _frames():
                loaded = []
                for p in images:
                    try:
                        loaded.append(np.array(Image.open(p).convert("RGB").resize((w, h), Image.LANCZOS)))
                    except Exception:
                        loaded.append(np.zeros((h, w, 3), dtype=np.uint8))
                for arr in loaded:
                    for i in range(tf):
                        yield (arr * (i / tf)).astype(np.uint8)
                    for _ in range(max(1, sf - tf * 2)):
                        yield arr
                    for i in range(tf):
                        yield (arr * (1.0 - i / tf)).astype(np.uint8)

            _write_frames(_frames(), out, sl_fps, w, h)
            return f"✅ 投影片影片已生成：{out}"

        # ── 文字動畫 ──────────────────────────────────────────────────
        elif mode == "text_video":
            text     = kwargs.get("text", "Hello")
            dur      = kwargs.get("duration", 5)
            bg_col   = tuple(kwargs.get("bg_color",   [30, 30, 40]))
            fg_col   = tuple(kwargs.get("font_color", [255, 255, 255]))
            fsize    = kwargs.get("font_size", 60)
            font     = _get_font(fsize)
            total    = max(1, int(fps * dur))

            def _frames():
                for i in range(total):
                    p = i / total
                    a = p / 0.2 if p < 0.2 else (1.0 - p) / 0.2 if p > 0.8 else 1.0
                    img  = Image.new("RGB", (w, h), bg_col)
                    draw = ImageDraw.Draw(img)
                    c    = tuple(int(x * a) for x in fg_col)
                    lines = text.split("\n")
                    lh    = fsize + 10
                    y0    = (h - len(lines) * lh) // 2
                    for j, ln in enumerate(lines):
                        bb = draw.textbbox((0, 0), ln, font=font)
                        draw.text(((w - (bb[2] - bb[0])) // 2, y0 + j * lh), ln, font=font, fill=c)
                    yield np.array(img)

            _write_frames(_frames(), out, fps, w, h)
            return f"✅ 文字動畫影片已生成：{out}"

        # ── TTS 語音影片 ─────────────────────────────────────────────
        elif mode == "tts_video":
            text      = kwargs.get("text", "")
            img_path  = kwargs.get("image", "")
            voice     = kwargs.get("voice", "zh-CN-YunxiNeural")
            subtitle  = kwargs.get("subtitle", True)
            if not text:
                return "❌ 請提供 text 參數"

            tmp     = Path(tempfile.mkdtemp())
            audio   = str(tmp / "tts.mp3")
            vidtmp  = str(tmp / "silent.mp4")

            _tts_sync(clean_for_tts(text), voice, audio)
            dur   = _get_audio_dur(audio)
            total = max(1, int(fps * dur))

            if img_path and Path(img_path).exists():
                bg = np.array(Image.open(img_path).convert("RGB").resize((w, h), Image.LANCZOS))
            else:
                bg = np.zeros((h, w, 3), dtype=np.uint8)
                for row in range(h):
                    v = int(20 + 30 * row / h)
                    bg[row, :] = [v, v, v + 20]

            font  = _get_font(40)
            cpl   = 20
            subs  = [text[i:i+cpl] for i in range(0, len(text), cpl)]

            def _frames():
                for i in range(total):
                    frame = bg.copy()
                    if subtitle and subs:
                        ln  = subs[min(int(i / total * len(subs)), len(subs) - 1)]
                        img = Image.fromarray(frame)
                        d   = ImageDraw.Draw(img)
                        bb  = d.textbbox((0, 0), ln, font=font)
                        tw  = bb[2] - bb[0]
                        tx  = (w - tw) // 2
                        ty  = int(h * 0.82)
                        d.rectangle([tx - 10, ty - 5, tx + tw + 10, ty + 45], fill=(0, 0, 0))
                        d.text((tx, ty), ln, font=font, fill=(255, 255, 255))
                        frame = np.array(img)
                    yield frame

            _write_frames(_frames(), vidtmp, fps, w, h)

            r = subprocess.run([
                ffmpeg_exe, "-y", "-i", vidtmp, "-i", audio,
                "-c:v", "copy", "-c:a", "aac", "-shortest", out
            ], capture_output=True)
            if not Path(out).exists():
                raise RuntimeError(f"音訊合併失敗：{r.stderr.decode(errors='replace')[-200:]}")

            return f"✅ TTS 語音影片已生成：{out}"

        # ── 螢幕錄影 ─────────────────────────────────────────────────
        elif mode == "screen_record":
            import mss, time as _t
            dur      = kwargs.get("duration", 10)
            rec_fps  = kwargs.get("fps", 10)
            interval = 1.0 / rec_fps
            total    = max(1, int(rec_fps * dur))

            with mss.mss() as sct:
                mon    = sct.monitors[1]
                sw, sh = mon["width"], mon["height"]

                def _frames():
                    for _ in range(total):
                        t0  = _t.time()
                        arr = np.array(sct.grab(mon))[:, :, :3][:, :, ::-1]
                        yield arr
                        elapsed = _t.time() - t0
                        if elapsed < interval:
                            _t.sleep(interval - elapsed)

                _write_frames(_frames(), out, rec_fps, sw, sh)

            return f"✅ 螢幕錄影完成：{out}（{dur} 秒）"

        else:
            return f"❌ 未知 mode：{mode}"

    except Exception as e:
        return f"❌ 影片生成失敗：{e}\n{traceback.format_exc()}"


def execute_monitor_config():
    try:
        from screeninfo import get_monitors
        return "\n".join(
            f"{'主螢幕' if m.is_primary else '副螢幕'} {m.width}x{m.height} @({m.x},{m.y}) {m.name}"
            for m in get_monitors()
        )
    except Exception as e:
        return f"❌ 取得螢幕資訊失敗：{e}"


def execute_run_code(type_, code):
    try:
        if type_ == "python":
            import io as _io, traceback, contextlib
            buf = _io.StringIO()
            with contextlib.redirect_stdout(buf):
                try:
                    exec(compile(code, "<string>", "exec"), {})
                except Exception:
                    buf.write(traceback.format_exc())
            return buf.getvalue() or "（執行完畢，無輸出）"
        elif type_ == "shell":
            result = subprocess.run(
                ["powershell.exe", "-Command", code],
                capture_output=True, text=True, encoding="utf-8", errors="replace", timeout=30
            )
            return (result.stdout + result.stderr).strip() or "（執行完畢，無輸出）"
        return f"不支援的程式碼類型「{type_}」，目前只支援 python 和 shell"
    except Exception as e:
        return f"執行失敗：{e}"


def execute_document(action, path, content="", sheet=None):
    try:
        if action == "word_read":
            from docx import Document
            doc = Document(path)
            return "\n".join(p.text for p in doc.paragraphs if p.text.strip()) or "（文件為空）"
        elif action == "word_write":
            from docx import Document
            doc = Document() if not Path(path).exists() else Document(path)
            doc.add_paragraph(content)
            doc.save(path)
            return f"已寫入：{path}"
        elif action == "excel_read":
            import openpyxl
            wb = openpyxl.load_workbook(path)
            ws = wb[sheet] if sheet else wb.active
            rows = [[str(c.value or "") for c in row] for row in ws.iter_rows()]
            return "\n".join(["\t".join(r) for r in rows[:30]])
        elif action == "excel_write":
            import openpyxl, json
            wb = openpyxl.load_workbook(path) if Path(path).exists() else openpyxl.Workbook()
            ws = wb[sheet] if (sheet and sheet in wb.sheetnames) else wb.active
            data = json.loads(content)
            for row in data:
                ws.append(row)
            wb.save(path)
            return f"已寫入 {len(data)} 行到 {path}"
        elif action == "pdf_read":
            import fitz
            doc = fitz.open(path)
            text = "\n".join(page.get_text() for page in doc)
            return text[:3000] if text else "（PDF 無可讀文字）"
        return "未知動作"
    except Exception as e:
        return f"文件操作失敗：{e}"


def execute_web_scrape(action, url="", selector="body", interval=2.0, region="full"):
    try:
        if action == "scrape":
            from bs4 import BeautifulSoup
            # 先用靜態爬蟲試，若內容太少（JS渲染網站）改用 Playwright
            res = requests.get(url, timeout=15, headers={"User-Agent": "Mozilla/5.0"})
            soup = BeautifulSoup(res.text, "html.parser")
            static_text = soup.get_text(strip=True)
            use_playwright = len(static_text) < 500 or len(soup.find_all("script")) > 10

            if use_playwright:
                try:
                    import subprocess, sys, json as _json, textwrap
                    _script = textwrap.dedent(f"""
import sys
from playwright.sync_api import sync_playwright
with sync_playwright() as p:
    b = p.chromium.launch(headless=True)
    page = b.new_page(user_agent="Mozilla/5.0 (Windows NT 10.0; Win64; x64) AppleWebKit/537.36")
    page.goto({_json.dumps(url)}, wait_until="networkidle", timeout=20000)
    page.wait_for_timeout(2000)
    els = page.query_selector_all({_json.dumps(selector)})
    texts = [e.inner_text() for e in els[:10] if e.inner_text().strip()]
    b.close()
    print("\\n".join(texts))
""")
                    proc = subprocess.run(
                        [sys.executable, "-c", _script],
                        capture_output=True, encoding="utf-8", timeout=30,
                        env={**__import__("os").environ, "PYTHONIOENCODING": "utf-8"}
                    )
                    if proc.returncode == 0 and proc.stdout.strip():
                        return proc.stdout.strip()
                    # Playwright 失敗 → fallback 回靜態內容
                except Exception:
                    pass

            # 靜態內容或 Playwright fallback
            elements = soup.select(selector)
            static_result = "\n".join(e.get_text(strip=True) for e in elements[:10])
            return static_result if static_result else f"（網頁內容太少，可能為 JS 動態網站：{url}）"
        elif action == "screen_diff":
            import time as t
            img1 = pyautogui.screenshot()
            t.sleep(interval)
            img2 = pyautogui.screenshot()
            import numpy as np
            a1, a2 = np.array(img1), np.array(img2)
            diff = np.abs(a1.astype(int) - a2.astype(int)).mean()
            if diff > 2.0:
                return f"⚠️ 螢幕有變化（差異度：{diff:.1f}）"
            return f"✅ 螢幕無明顯變化（差異度：{diff:.1f}）"
        return "未知動作"
    except Exception as e:
        return f"操作失敗：{e}"


def execute_image_edit(action, path, *params):
    try:
        from PIL import Image, ImageDraw, ImageFont
        img = Image.open(path)
        if action == "crop":
            x, y, w, h = int(params[0]), int(params[1]), int(params[2]), int(params[3])
            img = img.crop((x, y, x+w, y+h))
        elif action == "resize":
            w, h = int(params[0]), int(params[1])
            img = img.resize((w, h))
        elif action == "text":
            x, y = int(params[0]), int(params[1])
            text = " ".join(params[2:]) if len(params) > 2 else ""
            draw = ImageDraw.Draw(img)
            draw.text((x, y), text, fill="red")
        elif action == "merge":
            img2 = Image.open(params[0])
            merged = Image.new("RGB", (img.width + img2.width, max(img.height, img2.height)))
            merged.paste(img, (0, 0))
            merged.paste(img2, (img.width, 0))
            img = merged
        out = path.replace(".", f"_{action}.")
        img.save(out)
        return f"✅ 圖片已儲存：{out}"
    except Exception as e:
        return f"圖片編輯失敗：{e}"


def execute_cloud_storage(action, path, drive_id="root"):
    try:
        from google.oauth2.credentials import Credentials
        from googleapiclient.discovery import build
        from googleapiclient.http import MediaFileUpload, MediaIoBaseDownload
        import io as _io
        creds_path = Path("C:/Users/blue_/claude-telegram-bot/gdrive_token.json")
        if not creds_path.exists():
            return "❌ 未找到 Google Drive 憑證（gdrive_token.json）"
        creds = Credentials.from_authorized_user_file(str(creds_path))
        service = build("drive", "v3", credentials=creds)
        if action == "upload":
            media = MediaFileUpload(path)
            meta = {"name": Path(path).name, "parents": [drive_id]}
            f = service.files().create(body=meta, media_body=media, fields="id").execute()
            return f"✅ 已上傳，檔案 ID：{f.get('id')}"
        elif action == "download":
            req = service.files().get_media(fileId=drive_id)
            buf = _io.BytesIO()
            dl = MediaIoBaseDownload(buf, req)
            done = False
            while not done:
                _, done = dl.next_chunk()
            with open(path, "wb") as f:
                f.write(buf.getvalue())
            return f"✅ 已下載到：{path}"
        return "未知動作"
    except Exception as e:
        return f"雲端儲存失敗：{e}"


def execute_database(type_, db, sql, name=""):
    try:
        if type_ == "sqlite":
            conn = sqlite3.connect(db)
            cur = conn.cursor()
            cur.execute(sql)
            rows = cur.fetchall()
            conn.commit()
            conn.close()
            if rows:
                return "\n".join(str(r) for r in rows[:20])
            return "✅ 執行成功（無回傳資料）"
        elif type_ == "mysql":
            import pymysql
            conn = pymysql.connect(host=db, database=name, read_default_file="~/.my.cnf")
            cur = conn.cursor()
            cur.execute(sql)
            rows = cur.fetchall()
            conn.commit()
            conn.close()
            if rows:
                return "\n".join(str(r) for r in rows[:20])
            return "✅ 執行成功（無回傳資料）"
        return "未知類型"
    except Exception as e:
        return f"資料庫操作失敗：{e}"


def execute_encrypt_file(action, path, password):
    try:
        from cryptography.hazmat.primitives.kdf.pbkdf2 import PBKDF2HMAC
        from cryptography.hazmat.primitives import hashes
        from cryptography.fernet import Fernet
        import base64
        salt = b"xiaoniuma_salt_v1"
        kdf = PBKDF2HMAC(algorithm=hashes.SHA256(), length=32, salt=salt, iterations=100000)
        key = base64.urlsafe_b64encode(kdf.derive(password.encode()))
        f = Fernet(key)
        data = Path(path).read_bytes()
        if action == "encrypt":
            enc = f.encrypt(data)
            out = path + ".enc"
            Path(out).write_bytes(enc)
            return f"✅ 已加密：{out}"
        elif action == "decrypt":
            dec = f.decrypt(data)
            out = path.replace(".enc", ".dec")
            Path(out).write_bytes(dec)
            return f"✅ 已解密：{out}"
        return "未知動作"
    except Exception as e:
        return f"加密/解密失敗：{e}"


def execute_qr_code(action, content="", path="", duration=30.0):
    try:
        if action == "qr_gen":
            import qrcode as _qr
            img = _qr.make(content)
            save_path = path or str(Path("C:/Users/blue_/Desktop/測試檔案") / "qrcode.png")
            img.save(save_path)
            return f"✅ QR Code 已生成：{save_path}"
        elif action == "qr_scan":
            from pyzbar.pyzbar import decode
            from PIL import Image
            img = Image.open(path) if path else pyautogui.screenshot()
            results = decode(img)
            if not results:
                return "❌ 未偵測到 QR Code"
            return "\n".join(f"掃描結果：{r.data.decode('utf-8')}" for r in results)
        elif action == "clipboard_watch":
            import pyperclip, time as t
            last = pyperclip.paste()
            changes = []
            start = t.time()
            while t.time() - start < duration:
                cur = pyperclip.paste()
                if cur != last:
                    changes.append(f"[{datetime.datetime.now().strftime('%H:%M:%S')}] {cur[:100]}")
                    last = cur
                t.sleep(0.5)
            return "\n".join(changes) if changes else f"監控 {duration} 秒內無剪貼簿變化"
        return "未知動作"
    except Exception as e:
        return f"操作失敗：{e}"


# ── 缺口1：觸發驅動 ──────────────────────────────────────────────

def execute_email_trigger(action, host="", user="", password="", filter_from="",
                          filter_subject="", duration=300, to="", subject="", body=""):
    try:
        if action == "send":
            import smtplib
            from email.mime.text import MIMEText
            smtp_host = host.replace("imap.", "smtp.") if host else "smtp.gmail.com"
            msg = MIMEText(body, "plain", "utf-8")
            msg["Subject"] = subject
            msg["From"] = user
            msg["To"] = to
            with smtplib.SMTP_SSL(smtp_host, 465) as s:
                s.login(user, password)
                s.send_message(msg)
            return f"✅ 郵件已發送至 {to}"

        import imaplib, email as _email, time as _t
        def _connect():
            m = imaplib.IMAP4_SSL(host or "imap.gmail.com")
            m.login(user, password)
            m.select("INBOX")
            return m

        def _fetch_recent(m, n=5):
            _, data = m.search(None, "ALL")
            ids = data[0].split()[-n:]
            mails = []
            for mid in reversed(ids):
                _, md = m.fetch(mid, "(RFC822)")
                msg = _email.message_from_bytes(md[0][1])
                sender = msg.get("From", "")
                subj = msg.get("Subject", "")
                if filter_from and filter_from.lower() not in sender.lower():
                    continue
                if filter_subject and filter_subject.lower() not in subj.lower():
                    continue
                body_text = ""
                if msg.is_multipart():
                    for part in msg.walk():
                        if part.get_content_type() == "text/plain":
                            body_text = part.get_payload(decode=True).decode("utf-8", errors="ignore")[:300]
                            break
                else:
                    body_text = msg.get_payload(decode=True).decode("utf-8", errors="ignore")[:300]
                mails.append(f"寄件人：{sender}\n主旨：{subj}\n內容：{body_text}")
            return mails

        if action == "check":
            m = _connect()
            mails = _fetch_recent(m, 5)
            m.logout()
            return "\n---\n".join(mails) if mails else "📭 收件箱無符合郵件"

        elif action == "watch":
            m = _connect()
            seen = set()
            _, data = m.search(None, "ALL")
            for mid in data[0].split():
                seen.add(mid)
            results = []
            start = _t.time()
            while _t.time() - start < duration:
                _t.sleep(10)
                try:
                    m.noop()
                    _, data = m.search(None, "ALL")
                    for mid in data[0].split():
                        if mid not in seen:
                            seen.add(mid)
                            _, md = m.fetch(mid, "(RFC822)")
                            msg = _email.message_from_bytes(md[0][1])
                            results.append(f"📬 新郵件：{msg.get('From')} | {msg.get('Subject')}")
                except Exception:
                    m = _connect()
            m.logout()
            return "\n".join(results) if results else f"監控 {duration} 秒內無新郵件"
        return "未知動作"
    except Exception as e:
        return f"❌ email_trigger 失敗：{e}"


def execute_file_trigger(folder, event, action, pattern="", target="", duration=60):
    try:
        import time as _t, fnmatch, shutil, subprocess as _sp
        from pathlib import Path as _P
        folder_path = _P(folder)
        if not folder_path.exists():
            return f"❌ 資料夾不存在：{folder}"

        if action == "list":
            files = list(folder_path.iterdir())
            if pattern:
                files = [f for f in files if fnmatch.fnmatch(f.name, pattern)]
            return "\n".join(str(f) for f in files[:30]) or "資料夾為空"

        before = set(str(f) for f in folder_path.rglob("*"))
        triggered = []
        start = _t.time()
        while _t.time() - start < duration:
            _t.sleep(2)
            after = set(str(f) for f in folder_path.rglob("*"))
            if event in ("created", "any"):
                new_files = after - before
                for f in new_files:
                    if pattern and not fnmatch.fnmatch(_P(f).name, pattern):
                        continue
                    msg = f"[新增] {f}"
                    if action == "copy" and target:
                        shutil.copy2(f, target)
                        msg += f" → 已複製到 {target}"
                    elif action == "run" and target:
                        _sp.Popen([target, f])
                        msg += f" → 已執行 {target}"
                    elif action == "notify":
                        msg += " → 通知觸發"
                    triggered.append(msg)
            if event in ("deleted", "any"):
                removed = before - after
                for f in removed:
                    triggered.append(f"[刪除] {f}")
            before = after
        return "\n".join(triggered) if triggered else f"監控 {duration} 秒內無 {event} 事件"
    except Exception as e:
        return f"❌ file_trigger 失敗：{e}"


_webhook_server_proc = None

def execute_webhook_server(action, port=8765, secret=""):
    global _webhook_server_proc
    try:
        import socket, subprocess as _sp, sys, os as _os
        if action == "start":
            if _webhook_server_proc and _webhook_server_proc.poll() is None:
                return f"✅ Webhook 伺服器已在運行（port {port}）"
            script = f"""
import http.server, json, threading, time
log = []
class H(http.server.BaseHTTPRequestHandler):
    def do_POST(self):
        length = int(self.headers.get('Content-Length', 0))
        body = self.rfile.read(length).decode('utf-8', errors='ignore')
        secret = "{secret}"
        if secret and self.headers.get('X-Secret','') != secret:
            self.send_response(403); self.end_headers(); return
        log.append(f"[{{time.strftime('%H:%M:%S')}}] {{self.path}}: {{body[:200]}}")
        open(r'C:/Users/blue_/claude-telegram-bot/webhook_log.txt','a',encoding='utf-8').write(log[-1]+'\\n')
        self.send_response(200); self.send_header('Content-Type','text/plain'); self.end_headers()
        self.wfile.write(b'ok')
    def log_message(self, *a): pass
http.server.HTTPServer(('0.0.0.0', {port}), H).serve_forever()
"""
            _webhook_server_proc = _sp.Popen([sys.executable, "-c", script],
                                              creationflags=0x00000008)
            return f"✅ Webhook 伺服器已啟動 port {port}\n本機 IP：{socket.gethostbyname(socket.gethostname())}"
        elif action == "stop":
            if _webhook_server_proc:
                _webhook_server_proc.terminate()
                _webhook_server_proc = None
            return "✅ Webhook 伺服器已停止"
        elif action == "status":
            running = _webhook_server_proc and _webhook_server_proc.poll() is None
            log_path = r"C:/Users/blue_/claude-telegram-bot/webhook_log.txt"
            log = ""
            if _os.path.exists(log_path):
                with open(log_path, encoding="utf-8") as f:
                    log = "".join(f.readlines()[-5:])
            return f"狀態：{'運行中' if running else '已停止'}\n最近事件：\n{log}"
        elif action == "get_url":
            import socket
            ip = socket.gethostbyname(socket.gethostname())
            return f"Webhook URL：http://{ip}:{port}/\n（需同一區域網路）"
        return "未知動作"
    except Exception as e:
        return f"❌ webhook_server 失敗：{e}"


# ── 缺口2：應用程式深度控制 ──────────────────────────────────────

def execute_com_auto(app, action, path="", sheet=None, cell="", value="",
                     macro="", to="", subject=""):
    try:
        import win32com.client as _com
        if app == "excel":
            xl = _com.Dispatch("Excel.Application")
            xl.Visible = False
            wb = xl.Workbooks.Open(path) if path else (xl.Workbooks(1) if xl.Workbooks.Count > 0 else xl.Workbooks.Add())
            ws = wb.Sheets(sheet) if sheet else wb.ActiveSheet
            if action == "read_cell":
                return f"{cell} = {ws.Range(cell).Value}"
            elif action == "write_cell":
                ws.Range(cell).Value = value
                return f"✅ 已寫入 {cell} = {value}"
            elif action == "run_macro":
                xl.Run(macro)
                return f"✅ 巨集 {macro} 已執行"
            elif action == "save":
                wb.Save()
                return "✅ 已儲存"
            elif action == "list_sheets":
                return "\n".join(wb.Sheets(i+1).Name for i in range(wb.Sheets.Count))
            elif action == "close":
                wb.Close(SaveChanges=False)
                return "✅ 已關閉"
        elif app == "word":
            wd = _com.Dispatch("Word.Application")
            wd.Visible = False
            doc = wd.Documents.Open(path) if path else wd.Documents.Add()
            if action == "read":
                return doc.Content.Text[:1000]
            elif action == "write":
                doc.Content.Text = value
                return "✅ 已寫入"
            elif action == "save":
                doc.Save()
                return "✅ 已儲存"
            elif action == "close":
                doc.Close(SaveChanges=False)
                return "✅ 已關閉"
        elif app == "outlook":
            ol = _com.Dispatch("Outlook.Application")
            if action == "send":
                mail = ol.CreateItem(0)
                mail.To = to
                mail.Subject = subject
                mail.Body = value
                mail.Send()
                return f"✅ 郵件已發送至 {to}"
            elif action == "list_inbox":
                inbox = ol.GetNamespace("MAPI").GetDefaultFolder(6)
                items = inbox.Items
                items.Sort("[ReceivedTime]", True)
                result = []
                for i, item in enumerate(items):
                    if i >= 10:
                        break
                    result.append(f"寄件人：{item.SenderName} | 主旨：{item.Subject}")
                return "\n".join(result)
        return f"✅ {app} {action} 完成"
    except Exception as e:
        return f"❌ com_auto 失敗：{e}"


def execute_dialog_auto(action, button_text="", window_title="", timeout=30):
    try:
        import win32gui, win32con, time as _t
        def _find_dialogs():
            dialogs = []
            def cb(hwnd, _):
                if win32gui.IsWindowVisible(hwnd):
                    title = win32gui.GetWindowText(hwnd)
                    cls = win32gui.GetClassName(hwnd)
                    if cls in ("#32770", "TForm") or any(kw in title for kw in ["確認","警告","錯誤","提示","Dialog","Error","Warning","Confirm"]):
                        dialogs.append((hwnd, title, cls))
            win32gui.EnumWindows(cb, None)
            return dialogs

        def _click_button(hwnd, text):
            import win32api
            result = []
            def cb(child, _):
                t = win32gui.GetWindowText(child)
                if text.lower() in t.lower() or not text:
                    if win32gui.GetClassName(child) == "Button":
                        win32gui.SetForegroundWindow(hwnd)
                        win32api.SendMessage(child, win32con.BM_CLICK, 0, 0)
                        result.append(t)
            win32gui.EnumChildWindows(hwnd, cb, None)
            return result

        if action == "list_dialogs":
            dialogs = _find_dialogs()
            if not dialogs:
                return "目前無對話框"
            return "\n".join(f"HWND:{h} 標題:{t} 類別:{c}" for h, t, c in dialogs)

        elif action in ("find_and_click", "wait_and_click"):
            start = _t.time()
            btn = button_text or "確定"
            while True:
                dialogs = _find_dialogs()
                for hwnd, title, _ in dialogs:
                    if window_title and window_title.lower() not in title.lower():
                        continue
                    clicked = _click_button(hwnd, btn)
                    if clicked:
                        return f"✅ 已點擊對話框「{title}」的「{clicked[0]}」按鈕"
                if action == "find_and_click":
                    return "❌ 未找到符合的對話框"
                if _t.time() - start > timeout:
                    return f"❌ 等待 {timeout} 秒仍未出現對話框"
                _t.sleep(1)
        return "未知動作"
    except Exception as e:
        return f"❌ dialog_auto 失敗：{e}"


def execute_ime_switch(action):
    try:
        import ctypes, win32gui, win32api, win32con
        LANG_ZH = 0x0804
        LANG_EN = 0x0409

        hwnd = win32gui.GetForegroundWindow()
        if action == "status":
            lid = ctypes.windll.user32.GetKeyboardLayout(0) & 0xFFFF
            return f"目前輸入法 LCID：{hex(lid)} ({'中文' if lid in (0x0804,0x0404,0x0C04) else '英文' if lid == 0x0409 else '其他'})"
        elif action == "switch_en":
            hkl = win32api.LoadKeyboardLayout("00000409", 1)
            win32api.PostMessage(hwnd, win32con.WM_INPUTLANGCHANGEREQUEST, 0, hkl)
            return "✅ 已切換至英文輸入"
        elif action == "switch_zh":
            hkl = win32api.LoadKeyboardLayout("00000804", 1)
            win32api.PostMessage(hwnd, win32con.WM_INPUTLANGCHANGEREQUEST, 0, hkl)
            return "✅ 已切換至中文輸入"
        elif action == "toggle":
            import pyautogui
            pyautogui.hotkey("shift", "alt")
            return "✅ 已切換輸入法"
        return "未知動作"
    except Exception as e:
        return f"❌ ime_switch 失敗：{e}"


# ── 缺口3：感知能力 ──────────────────────────────────────────────

def execute_wake_word(action, keyword="", duration=5, language="zh-TW"):
    try:
        import speech_recognition as _sr
        r = _sr.Recognizer()
        mic = _sr.Microphone()

        if action == "listen_once":
            with mic as src:
                r.adjust_for_ambient_noise(src, duration=0.5)
                audio = r.listen(src, timeout=duration, phrase_time_limit=duration)
            try:
                text = r.recognize_google(audio, language=language)
                return f"🎤 聽到：{text}"
            except _sr.UnknownValueError:
                return "❌ 無法辨識語音"

        elif action == "transcribe_stream":
            import time as _t
            results = []
            end = _t.time() + duration
            while _t.time() < end:
                with mic as src:
                    r.adjust_for_ambient_noise(src, duration=0.3)
                    try:
                        audio = r.listen(src, timeout=2, phrase_time_limit=5)
                        text = r.recognize_google(audio, language=language)
                        results.append(text)
                    except (_sr.WaitTimeoutError, _sr.UnknownValueError):
                        pass
            return "\n".join(results) if results else "未偵測到語音"

        elif action == "detect_keyword":
            import time as _t
            end = _t.time() + duration
            while _t.time() < end:
                with mic as src:
                    r.adjust_for_ambient_noise(src, duration=0.3)
                    try:
                        audio = r.listen(src, timeout=2, phrase_time_limit=5)
                        text = r.recognize_google(audio, language=language)
                        if keyword.lower() in text.lower():
                            return f"✅ 偵測到關鍵字「{keyword}」，完整語音：{text}"
                    except (_sr.WaitTimeoutError, _sr.UnknownValueError):
                        pass
            return f"❌ {duration} 秒內未偵測到關鍵字「{keyword}」"
        return "未知動作"
    except ImportError:
        return "❌ 請先安裝：pip install SpeechRecognition"
    except Exception as e:
        return f"❌ wake_word 失敗：{e}"


def execute_sound_detect(action, threshold=20, duration=5, output=""):
    try:
        import numpy as _np
        try:
            import sounddevice as _sd
        except ImportError:
            return "❌ 請先安裝：pip install sounddevice"

        RATE = 16000
        if action == "volume_level":
            data = _sd.rec(int(RATE * 1), samplerate=RATE, channels=1, dtype="float32")
            _sd.wait()
            vol = int(_np.abs(data).mean() * 1000)
            return f"🔊 當前音量：{vol}/100"

        elif action == "detect_silence":
            import time as _t
            silent_start = None
            start = _t.time()
            while _t.time() - start < duration:
                data = _sd.rec(int(RATE * 0.5), samplerate=RATE, channels=1, dtype="float32")
                _sd.wait()
                vol = int(_np.abs(data).mean() * 1000)
                if vol < threshold:
                    if silent_start is None:
                        silent_start = _t.time()
                    elif _t.time() - silent_start > 1.5:
                        return f"🔇 偵測到靜音（音量 {vol}）"
                else:
                    silent_start = None
            return f"監控 {duration} 秒內無靜音段"

        elif action == "detect_speech":
            import time as _t
            start = _t.time()
            while _t.time() - start < duration:
                data = _sd.rec(int(RATE * 0.5), samplerate=RATE, channels=1, dtype="float32")
                _sd.wait()
                vol = int(_np.abs(data).mean() * 1000)
                if vol > threshold:
                    return f"🗣 偵測到說話（音量 {vol}）"
            return f"監控 {duration} 秒內未偵測到說話"

        elif action == "record_until_silence":
            import time as _t
            import wave
            frames = []
            out = output or str(__import__("pathlib").Path("C:/Users/blue_/Desktop/測試檔案") / "recording.wav")
            CHUNK = int(RATE * 0.5)
            silent_count = 0
            while True:
                data = _sd.rec(CHUNK, samplerate=RATE, channels=1, dtype="int16")
                _sd.wait()
                frames.append(data.tobytes())
                vol = int(_np.abs(data.astype(float)).mean() / 32768 * 1000)
                if vol < threshold:
                    silent_count += 1
                    if silent_count >= 4:
                        break
                else:
                    silent_count = 0
                if len(frames) > 200:
                    break
            with wave.open(out, "wb") as wf:
                wf.setnchannels(1)
                wf.setsampwidth(2)
                wf.setframerate(RATE)
                wf.writeframes(b"".join(frames))
            return f"✅ 錄音完成：{out}（{len(frames)*0.5:.1f}秒）"
        return "未知動作"
    except Exception as e:
        return f"❌ sound_detect 失敗：{e}"


def execute_face_recognize(action, name="", image_path="", output=""):
    try:
        import cv2 as _cv2
        import numpy as _np
        from pathlib import Path as _P
        import json, os as _os

        FACE_DB = str(_P.home() / ".face_db")
        _os.makedirs(FACE_DB, exist_ok=True)

        try:
            import face_recognition as _fr
        except ImportError:
            return "❌ 請先安裝：pip install face-recognition opencv-python"

        def _capture_frame():
            cap = _cv2.VideoCapture(0)
            ret, frame = cap.read()
            cap.release()
            if not ret:
                raise RuntimeError("無法開啟攝影機")
            return frame

        if action == "capture":
            frame = _capture_frame()
            out = output or str(_P("C:/Users/blue_/Desktop/測試檔案") / "face_capture.jpg")
            _cv2.imwrite(out, frame)
            return f"✅ 已拍照：{out}"

        elif action == "detect":
            img = _fr.load_image_file(image_path) if image_path else _capture_frame()
            if image_path:
                img = _cv2.cvtColor(_cv2.imread(image_path), _cv2.COLOR_BGR2RGB)
            locs = _fr.face_locations(img)
            return f"✅ 偵測到 {len(locs)} 個人臉"

        elif action == "enroll":
            if not name:
                return "❌ 需提供 name"
            img = _fr.load_image_file(image_path) if image_path else _capture_frame()
            if not image_path:
                img = _cv2.cvtColor(img, _cv2.COLOR_BGR2RGB)
            encs = _fr.face_encodings(img)
            if not encs:
                return "❌ 未偵測到人臉"
            enc_path = _os.path.join(FACE_DB, f"{name}.npy")
            _np.save(enc_path, encs[0])
            return f"✅ 已登記人臉：{name}"

        elif action == "recognize":
            known_encs, known_names = [], []
            for f in _P(FACE_DB).glob("*.npy"):
                known_encs.append(_np.load(str(f)))
                known_names.append(f.stem)
            if not known_encs:
                return "❌ 尚未登記任何人臉，請先用 enroll"
            img = _fr.load_image_file(image_path) if image_path else _capture_frame()
            if not image_path:
                img = _cv2.cvtColor(img, _cv2.COLOR_BGR2RGB)
            encs = _fr.face_encodings(img)
            if not encs:
                return "❌ 未偵測到人臉"
            results = []
            for enc in encs:
                matches = _fr.compare_faces(known_encs, enc)
                dists = _fr.face_distance(known_encs, enc)
                best = int(_np.argmin(dists))
                if matches[best]:
                    results.append(f"✅ 識別為：{known_names[best]}（相似度 {(1-dists[best])*100:.0f}%）")
                else:
                    results.append("❓ 未知人物")
            return "\n".join(results)
        return "未知動作"
    except Exception as e:
        return f"❌ face_recognize 失敗：{e}"


# ── 缺口4：跨裝置控制 ────────────────────────────────────────────

_http_control_proc = None

def execute_http_server(action, port=9876, password=""):
    global _http_control_proc
    try:
        import socket, subprocess as _sp, sys, os as _os
        if action == "start":
            if _http_control_proc and _http_control_proc.poll() is None:
                ip = socket.gethostbyname(socket.gethostname())
                return f"✅ HTTP 控制伺服器已在運行\n網址：http://{ip}:{port}/"
            script = f"""
import http.server, urllib.parse, json, subprocess, sys, os
PASSWORD = "{password}"
class H(http.server.BaseHTTPRequestHandler):
    def do_GET(self):
        if PASSWORD and self.headers.get('X-Password','') != PASSWORD:
            p = urllib.parse.urlparse(self.path)
            q = urllib.parse.parse_qs(p.query)
            if q.get('pw',[''])[0] != PASSWORD:
                self.send_response(403); self.end_headers()
                self.wfile.write(b'Unauthorized'); return
        self.send_response(200)
        self.send_header('Content-Type','text/html; charset=utf-8')
        self.end_headers()
        html = '''<html><body><h2>小牛馬遠端控制</h2>
<form method=post><input name=cmd size=60 placeholder="輸入指令"><button>執行</button></form></body></html>'''
        self.wfile.write(html.encode())
    def do_POST(self):
        l = int(self.headers.get('Content-Length',0))
        body = urllib.parse.parse_qs(self.rfile.read(l).decode())
        cmd = body.get('cmd',[''])[0]
        self.send_response(200)
        self.send_header('Content-Type','text/plain; charset=utf-8')
        self.end_headers()
        if cmd:
            try:
                out = subprocess.check_output(cmd, shell=True, text=True, stderr=subprocess.STDOUT, timeout=10)
            except Exception as e:
                out = str(e)
            self.wfile.write(out.encode('utf-8','ignore'))
    def log_message(self, *a): pass
http.server.HTTPServer(('0.0.0.0', {port}), H).serve_forever()
"""
            _http_control_proc = _sp.Popen([sys.executable, "-c", script],
                                            creationflags=0x00000008)
            ip = socket.gethostbyname(socket.gethostname())
            return f"✅ HTTP 控制伺服器已啟動\n網址：http://{ip}:{port}/\n（手機在同一 WiFi 下可存取）"
        elif action == "stop":
            if _http_control_proc:
                _http_control_proc.terminate()
                _http_control_proc = None
            return "✅ 已停止"
        elif action == "status":
            running = _http_control_proc and _http_control_proc.poll() is None
            return f"狀態：{'運行中' if running else '已停止'}"
        elif action == "get_url":
            ip = socket.gethostbyname(socket.gethostname())
            return f"http://{ip}:{port}/"
        return "未知動作"
    except Exception as e:
        return f"❌ http_server 失敗：{e}"


def execute_lan_scan(action, subnet="", host="", port=80):
    try:
        import socket, subprocess as _sp, concurrent.futures
        if action == "get_local_ip":
            s = socket.socket(socket.AF_INET, socket.SOCK_DGRAM)
            s.connect(("8.8.8.8", 80))
            ip = s.getsockname()[0]
            s.close()
            return f"本機 IP：{ip}"

        elif action == "ping_sweep":
            if not subnet:
                s = socket.socket(socket.AF_INET, socket.SOCK_DGRAM)
                s.connect(("8.8.8.8", 80))
                local_ip = s.getsockname()[0]
                s.close()
                subnet = ".".join(local_ip.split(".")[:3]) + ".0/24"
            base = ".".join(subnet.split(".")[:3])
            online = []
            def ping(i):
                ip = f"{base}.{i}"
                r = _sp.run(["ping", "-n", "1", "-w", "300", ip],
                            capture_output=True, text=True)
                if "TTL=" in r.stdout or "ttl=" in r.stdout:
                    try:
                        hostname = socket.gethostbyaddr(ip)[0]
                    except Exception:
                        hostname = "unknown"
                    return f"{ip} ({hostname})"
                return None
            with concurrent.futures.ThreadPoolExecutor(max_workers=50) as ex:
                results = ex.map(ping, range(1, 255))
            online = [r for r in results if r]
            return f"區域網路掃描結果（{subnet}）：\n" + "\n".join(online) if online else "無裝置回應"

        elif action == "port_check":
            s = socket.socket(socket.AF_INET, socket.SOCK_STREAM)
            s.settimeout(2)
            result = s.connect_ex((host, int(port)))
            s.close()
            return f"{'✅' if result == 0 else '❌'} {host}:{port} {'開啟' if result == 0 else '關閉'}"

        elif action == "scan":
            s = socket.socket(socket.AF_INET, socket.SOCK_DGRAM)
            s.connect(("8.8.8.8", 80))
            local_ip = s.getsockname()[0]
            s.close()
            base = ".".join(local_ip.split(".")[:3])
            common_ports = [21, 22, 80, 443, 3389, 8080, 1883]
            results = []
            def check(i):
                ip = f"{base}.{i}"
                r = _sp.run(["ping", "-n", "1", "-w", "200", ip], capture_output=True)
                if b"TTL=" in r.stdout or b"ttl=" in r.stdout:
                    open_ports = []
                    for p in common_ports:
                        sock = socket.socket()
                        sock.settimeout(0.3)
                        if sock.connect_ex((ip, p)) == 0:
                            open_ports.append(str(p))
                        sock.close()
                    try:
                        hn = socket.gethostbyaddr(ip)[0]
                    except Exception:
                        hn = "?"
                    return f"{ip} ({hn}) 開放 port: {','.join(open_ports) or '無'}"
                return None
            with concurrent.futures.ThreadPoolExecutor(max_workers=50) as ex:
                res = list(ex.map(check, range(1, 255)))
            found = [r for r in res if r]
            return "\n".join(found) if found else "未找到裝置"
        return "未知動作"
    except Exception as e:
        return f"❌ lan_scan 失敗：{e}"


def execute_serial_port(action, port="", baudrate=9600, data="", timeout=2):
    try:
        import serial, serial.tools.list_ports
        if action == "list":
            ports = serial.tools.list_ports.comports()
            if not ports:
                return "❌ 未找到 COM port"
            return "\n".join(f"{p.device} - {p.description}" for p in ports)
        if not port:
            return "❌ 需指定 port"
        if action == "send":
            with serial.Serial(port, baudrate, timeout=timeout) as s:
                s.write(data.encode())
            return f"✅ 已發送：{data}"
        elif action == "read":
            with serial.Serial(port, baudrate, timeout=timeout) as s:
                resp = s.read(1024).decode("utf-8", errors="ignore")
            return f"收到：{resp}" if resp else "❌ 無資料回應"
        elif action == "send_read":
            with serial.Serial(port, baudrate, timeout=timeout) as s:
                s.write(data.encode())
                import time; time.sleep(0.1)
                resp = s.read(1024).decode("utf-8", errors="ignore")
            return f"發送：{data}\n收到：{resp}"
        return "未知動作"
    except ImportError:
        return "❌ 請先安裝：pip install pyserial"
    except Exception as e:
        return f"❌ serial_port 失敗：{e}"


def execute_mqtt(action, broker, port=1883, topic="", message="",
                 duration=10, username="", password=""):
    try:
        import paho.mqtt.client as _mqtt
        import time as _t
    except ImportError:
        return "❌ 請先安裝：pip install paho-mqtt"
    try:
        received = []
        connected = [False]

        def on_connect(c, ud, flags, rc):
            connected[0] = rc == 0

        def on_message(c, ud, msg):
            received.append(f"[{msg.topic}] {msg.payload.decode('utf-8','ignore')}")

        client = _mqtt.Client()
        client.on_connect = on_connect
        client.on_message = on_message
        if username:
            client.username_pw_set(username, password)
        client.connect(broker, int(port), 60)
        client.loop_start()
        _t.sleep(1)

        if not connected[0]:
            return f"❌ 無法連線至 {broker}:{port}"

        if action == "test_connect":
            client.disconnect()
            return f"✅ 成功連線 {broker}:{port}"
        elif action == "publish":
            client.publish(topic, message)
            _t.sleep(0.5)
            client.disconnect()
            return f"✅ 已發布到 {topic}：{message}"
        elif action == "subscribe":
            client.subscribe(topic)
            _t.sleep(duration)
            client.disconnect()
            return "\n".join(received) if received else f"訂閱 {topic} 共 {duration} 秒，無訊息"
        return "未知動作"
    except Exception as e:
        return f"❌ mqtt 失敗：{e}"


# ── 缺口5：內容理解與處理 ────────────────────────────────────────

def execute_doc_ai(action, path="", path2="", fields="", question="", url=""):
    try:
        import anthropic as _ant, base64, mimetypes
        from pathlib import Path as _P

        client = _ant.Anthropic(api_key=os.getenv("ANTHROPIC_API_KEY", ""))

        def _load_content(p, u=""):
            if u:
                return [{"type": "text", "text": f"請分析這個網址的內容：{u}"}]
            if not p or not _P(p).exists():
                return [{"type": "text", "text": f"（檔案不存在：{p}）"}]
            ext = _P(p).suffix.lower()
            if ext in (".png", ".jpg", ".jpeg", ".gif", ".webp"):
                with open(p, "rb") as f:
                    data = base64.b64encode(f.read()).decode()
                mime = mimetypes.guess_type(p)[0] or "image/jpeg"
                return [{"type": "image", "source": {"type": "base64", "media_type": mime, "data": data}}]
            elif ext == ".pdf":
                try:
                    import pdfplumber
                    with pdfplumber.open(p) as pdf:
                        text = "\n".join(pg.extract_text() or "" for pg in pdf.pages[:10])
                except ImportError:
                    text = f"（需安裝 pdfplumber 才能讀取 PDF：{p}）"
                return [{"type": "text", "text": text[:4000]}]
            else:
                with open(p, "r", encoding="utf-8", errors="ignore") as f:
                    return [{"type": "text", "text": f.read()[:4000]}]

        content = _load_content(path, url)

        if action == "extract":
            prompt = f"請從以下內容提取這些欄位：{fields or '所有重要資訊'}。以 JSON 格式回傳。"
        elif action == "summarize":
            prompt = "請用繁體中文摘要這份文件的主要內容（200字以內）。"
        elif action == "classify":
            prompt = "請判斷這份文件的類型（如：發票、合約、報表、履歷等），並說明判斷依據。"
        elif action == "qa":
            prompt = f"根據文件內容回答：{question}"
        elif action == "compare":
            content2 = _load_content(path2)
            prompt = "請比較這兩份文件的差異，列出主要不同點。"
            content = content + [{"type": "text", "text": "---第二份文件---"}] + content2
        else:
            prompt = "請分析並說明這份文件的內容。"

        resp = client.messages.create(
            model="claude-sonnet-4-6",
            max_tokens=1024,
            messages=[{"role": "user", "content": content + [{"type": "text", "text": prompt}]}]
        )
        return resp.content[0].text
    except Exception as e:
        return f"❌ doc_ai 失敗：{e}"


def execute_web_monitor(action, url, selector="body", interval=60, duration=300, keyword=""):
    try:
        import requests as _req, time as _t, hashlib, json
        from pathlib import Path as _P
        try:
            from bs4 import BeautifulSoup as _BS
        except ImportError:
            _BS = None

        CACHE_FILE = str(_P.home() / ".web_monitor_cache.json")

        def _fetch(u, sel):
            r = _req.get(u, timeout=15, headers={"User-Agent": "Mozilla/5.0"})
            r.raise_for_status()
            if _BS:
                soup = _BS(r.text, "html.parser")
                el = soup.select_one(sel)
                return el.get_text(strip=True)[:2000] if el else soup.get_text(strip=True)[:2000]
            return r.text[:2000]

        if action == "check_once":
            text = _fetch(url, selector)
            return f"📄 {url}\n{text}"

        elif action == "get_price":
            text = _fetch(url, selector)
            import re
            prices = re.findall(r"[\$NT\$￥]?\s*[\d,]+\.?\d*", text)
            return f"💰 找到的價格：{', '.join(prices[:10])}" if prices else "❌ 未找到價格"

        elif action == "diff":
            try:
                with open(CACHE_FILE, encoding="utf-8") as f:
                    cache = json.load(f)
            except Exception:
                cache = {}
            current = _fetch(url, selector)
            h = hashlib.md5(current.encode()).hexdigest()
            prev = cache.get(url, "")
            cache[url] = current
            with open(CACHE_FILE, "w", encoding="utf-8") as f:
                json.dump(cache, f, ensure_ascii=False)
            if not prev:
                return f"✅ 已記錄初始狀態（{len(current)} 字）"
            if prev == current:
                return "✅ 與上次相同，無變化"
            import difflib
            diff = list(difflib.unified_diff(prev.splitlines(), current.splitlines(), lineterm="", n=2))
            return "⚠️ 內容有變化：\n" + "\n".join(diff[:30])

        elif action == "watch":
            last = _fetch(url, selector)
            changes = []
            start = _t.time()
            while _t.time() - start < duration:
                _t.sleep(interval)
                try:
                    current = _fetch(url, selector)
                    if current != last:
                        msg = f"⚠️ [{_t.strftime('%H:%M:%S')}] 網頁內容變化"
                        if keyword:
                            if keyword in current and keyword not in last:
                                msg += f"（出現關鍵字「{keyword}」）"
                            elif keyword not in current and keyword in last:
                                msg += f"（關鍵字「{keyword}」消失）"
                        changes.append(msg)
                        last = current
                except Exception as ex:
                    changes.append(f"❌ 抓取失敗：{ex}")
            return "\n".join(changes) if changes else f"監控 {duration} 秒內無變化"
        return "未知動作"
    except Exception as e:
        return f"❌ web_monitor 失敗：{e}"


def execute_audio_transcribe(action, path="", duration=30, language="", output=""):
    try:
        if action == "transcribe_file":
            if not path:
                return "❌ 需提供 path"
            try:
                import whisper
                model = whisper.load_model("base")
                result = model.transcribe(path, language=language or None)
                text = result["text"]
            except ImportError:
                import speech_recognition as _sr
                r = _sr.Recognizer()
                with _sr.AudioFile(path) as src:
                    audio = r.record(src)
                lang = language or "zh-TW"
                text = r.recognize_google(audio, language=lang)
            if output:
                with open(output, "w", encoding="utf-8") as f:
                    f.write(text)
                return f"✅ 轉錄完成，已儲存：{output}\n\n{text[:500]}"
            return f"📝 轉錄結果：\n{text}"

        elif action == "transcribe_mic":
            import speech_recognition as _sr
            r = _sr.Recognizer()
            m = _sr.Microphone()
            results = []
            import time as _t
            end = _t.time() + duration
            while _t.time() < end:
                with m as src:
                    r.adjust_for_ambient_noise(src, duration=0.3)
                    try:
                        audio = r.listen(src, timeout=3, phrase_time_limit=10)
                        lang = language or "zh-TW"
                        text = r.recognize_google(audio, language=lang)
                        results.append(text)
                    except (_sr.WaitTimeoutError, _sr.UnknownValueError):
                        pass
            full = " ".join(results)
            if output:
                with open(output, "w", encoding="utf-8") as f:
                    f.write(full)
            return f"📝 麥克風轉錄：\n{full}" if full else "❌ 未偵測到語音"

        elif action == "transcribe_system":
            return "⚠️ 系統音訊轉錄需要虛擬音訊裝置（如 VB-Cable），建議先錄製後用 transcribe_file"
        return "未知動作"
    except ImportError:
        return "❌ 請先安裝：pip install SpeechRecognition（或 openai-whisper 提供更高精度）"
    except Exception as e:
        return f"❌ audio_transcribe 失敗：{e}"


def fetch_image(prompt: str, width: int = 512, height: int = 512):
    hf_token = os.getenv("HF_TOKEN")
    headers = {"Authorization": f"Bearer {hf_token}"}
    payload = {"inputs": prompt}
    models = [
        "black-forest-labs/FLUX.1-schnell",
        "stabilityai/stable-diffusion-xl-base-1.0",
    ]
    for model in models:
        for attempt in range(2):
            try:
                res = requests.post(
                    f"https://router.huggingface.co/hf-inference/models/{model}",
                    headers=headers,
                    json=payload,
                    timeout=120
                )
                if res.status_code == 200 and res.headers.get("content-type", "").startswith("image"):
                    return res.content
                if res.status_code == 503:
                    time.sleep(10)
                    continue
            except Exception:
                pass
    return None


# ══════════════════════════════════════════════════════
# 奧創升級技能集 v1.0
# ══════════════════════════════════════════════════════

def execute_osint_search(action, query="", target="", limit=10):
    """OSINT情報蒐集：搜尋網路、新聞、社群媒體"""
    import feedparser
    results = []

    if action == "web_search":
        try:
            from ddgs import DDGS
            with DDGS() as ddgs:
                items = list(ddgs.text(query, region="zh-tw", max_results=int(limit)))
            for r in items:
                results.append({
                    "title": r.get("title", ""),
                    "url": r.get("href", ""),
                    "snippet": r.get("body", "")
                })
        except Exception as e:
            return f"搜尋失敗：{e}"
        return json.dumps(results, ensure_ascii=False, indent=2) if results else "無結果"

    elif action == "news_search":
        try:
            import feedparser
            url = f"https://news.google.com/rss/search?q={urllib.parse.quote(query)}&hl=zh-TW&gl=TW&ceid=TW:zh-Hant"
            feed = feedparser.parse(url)
            for entry in feed.entries[:int(limit)]:
                results.append(f"📰 {entry.get('title','')}\n   {entry.get('published','')}\n   {entry.get('link','')[:100]}")
        except Exception as e:
            return f"新聞搜尋失敗：{e}"
        return "\n\n".join(results) if results else "無相關新聞"

    elif action == "reddit_search":
        try:
            url = f"https://www.reddit.com/search.json?q={urllib.parse.quote(query)}&limit={limit}&sort=relevance"
            headers = {"User-Agent": "OsintBot/1.0"}
            resp = requests.get(url, headers=headers, timeout=10)
            data = resp.json()
            for post in data.get("data", {}).get("children", []):
                p = post["data"]
                results.append(f"r/{p.get('subreddit','')} | {p.get('title','')} (👍{p.get('score',0)})\n{p.get('url','')}")
        except Exception as e:
            return f"Reddit搜尋失敗：{e}"
        return "\n\n".join(results) if results else "無結果"

    elif action == "ip_osint":
        try:
            resp = requests.get(f"http://ip-api.com/json/{target}", timeout=5)
            d = resp.json()
            lines = [f"🌐 IP情報：{target}",
                     f"國家：{d.get('country','')} ({d.get('countryCode','')})",
                     f"城市：{d.get('city','')} / {d.get('regionName','')}",
                     f"ISP：{d.get('isp','')}",
                     f"組織：{d.get('org','')}",
                     f"時區：{d.get('timezone','')}",
                     f"座標：{d.get('lat','')}, {d.get('lon','')}"]
            return "\n".join(lines)
        except Exception as e:
            return f"IP查詢失敗：{e}"

    elif action == "domain_osint":
        try:
            import socket
            lines = [f"🔍 域名情報：{target}"]
            try:
                ip = socket.gethostbyname(target)
                lines.append(f"IP：{ip}")
                resp = requests.get(f"http://ip-api.com/json/{ip}", timeout=5)
                d = resp.json()
                lines.append(f"地區：{d.get('country','')} / {d.get('city','')}")
                lines.append(f"ISP：{d.get('isp','')}")
            except Exception as e2:
                lines.append(f"DNS解析失敗：{e2}")
            return "\n".join(lines)
        except Exception as e:
            return f"域名查詢失敗：{e}"

    elif action == "top_news":
        try:
            import feedparser
            feed = feedparser.parse("https://news.google.com/rss?hl=zh-TW&gl=TW&ceid=TW:zh-Hant")
            result = ["📡 Google 台灣頭條新聞："]
            for entry in feed.entries[:10]:
                result.append(f"• {entry.get('title','')}")
            return "\n".join(result)
        except Exception as e:
            return f"取得新聞失敗：{e}"

    return f"未知動作：{action}"


def execute_news_monitor(action, keywords="", interval=300, duration=3600, chat_id=None, _bot_send=None):
    """全球新聞監控：持續追蹤關鍵字新聞並即時通知"""
    import feedparser, threading

    if action == "check":
        results = []
        kw_list = [k.strip() for k in keywords.split(",") if k.strip()]
        for kw in kw_list[:5]:
            try:
                url = f"https://news.google.com/rss/search?q={urllib.parse.quote(kw)}&hl=zh-TW"
                feed = feedparser.parse(url)
                if feed.entries:
                    results.append(f"【{kw}】")
                    for entry in feed.entries[:3]:
                        results.append(f"  • {entry.get('title','')}")
            except Exception:
                pass
        return "\n".join(results) if results else "無相關新聞"

    elif action == "top_headlines":
        try:
            cats = [
                ("台灣", "https://news.google.com/rss?hl=zh-TW&gl=TW&ceid=TW:zh-Hant"),
                ("科技", "https://news.google.com/rss/topics/CAAqJggKIiBDQkFTRWdvSUwyMHZNRGRqTVhZU0FtZGhLQUFQAQ?hl=zh-TW"),
            ]
            result = []
            for cat, url in cats:
                feed = feedparser.parse(url)
                result.append(f"【{cat}】")
                for e in feed.entries[:4]:
                    result.append(f"  • {e.get('title','')}")
            return "\n".join(result)
        except Exception as e:
            return f"取得新聞失敗：{e}"

    elif action == "start_watch":
        seen = set()
        kw_list = [k.strip() for k in keywords.split(",") if k.strip()]
        start_t = time.time()

        def monitor():
            while time.time() - start_t < float(duration):
                for kw in kw_list:
                    try:
                        url = f"https://news.google.com/rss/search?q={urllib.parse.quote(kw)}&hl=zh-TW"
                        feed = feedparser.parse(url)
                        for entry in feed.entries[:3]:
                            link = entry.get("link", "")
                            if link and link not in seen:
                                seen.add(link)
                                msg = f"🚨 新聞快報 [{kw}]\n{entry.get('title','')}"
                                if _bot_send and chat_id:
                                    import asyncio
                                    asyncio.run_coroutine_threadsafe(
                                        _bot_send(chat_id=chat_id, text=msg),
                                        asyncio.get_event_loop()
                                    )
                    except Exception:
                        pass
                time.sleep(float(interval))

        threading.Thread(target=monitor, daemon=True).start()
        return f"✅ 開始監控：{keywords}，每 {interval}秒 檢查一次，持續 {int(float(duration)//3600)} 小時"

    return f"未知動作：{action}"


def execute_threat_intel(action, target="", api_key=""):
    """威脅情報查詢：IP/域名/Hash惡意程度分析"""

    if action in ("check_url", "check_hash", "check_ip"):
        vt_key = api_key or os.environ.get("VIRUSTOTAL_KEY", "")
        if not vt_key:
            return "需要設定環境變數 VIRUSTOTAL_KEY"
        headers = {"x-apikey": vt_key}
        try:
            if action == "check_url":
                import base64
                url_id = base64.urlsafe_b64encode(target.encode()).decode().strip("=")
                resp = requests.get(f"https://www.virustotal.com/api/v3/urls/{url_id}", headers=headers, timeout=10)
            elif action == "check_ip":
                resp = requests.get(f"https://www.virustotal.com/api/v3/ip_addresses/{target}", headers=headers, timeout=10)
            else:
                resp = requests.get(f"https://www.virustotal.com/api/v3/files/{target}", headers=headers, timeout=10)

            if resp.status_code == 200:
                stats = resp.json().get("data", {}).get("attributes", {}).get("last_analysis_stats", {})
                malicious = stats.get("malicious", 0)
                suspicious = stats.get("suspicious", 0)
                total = sum(stats.values()) if stats else 0
                level = "🔴 高危" if malicious > 3 else ("🟡 可疑" if malicious > 0 else "🟢 安全")
                return f"威脅分析：{target}\n狀態：{level}\n惡意：{malicious}/{total}\n可疑：{suspicious}/{total}"
            return f"查詢失敗：HTTP {resp.status_code}"
        except Exception as e:
            return f"查詢失敗：{e}"

    elif action == "scan_connections":
        try:
            import subprocess
            result = subprocess.run(["netstat", "-ano"], capture_output=True, text=True, timeout=10)
            lines = [l for l in result.stdout.split("\n") if "ESTABLISHED" in l]
            external = [l for l in lines if not any(ip in l for ip in ["127.0.0.1", "0.0.0.0", "::1", "[::1]"])]
            return f"活躍外部連線（{len(external)} 條）：\n" + "\n".join(external[:15])
        except Exception as e:
            return f"掃描失敗：{e}"

    elif action == "check_abuse_ip":
        key = api_key or os.environ.get("ABUSEIPDB_KEY", "")
        if not key:
            return "需要設定環境變數 ABUSEIPDB_KEY"
        try:
            resp = requests.get(
                "https://api.abuseipdb.com/api/v2/check",
                params={"ipAddress": target, "maxAgeInDays": 90},
                headers={"Key": key, "Accept": "application/json"}, timeout=10
            )
            if resp.status_code == 200:
                d = resp.json().get("data", {})
                return f"IP威脅：{target}\n濫用信心度：{d.get('abuseConfidenceScore',0)}%\n舉報次數：{d.get('totalReports',0)}\n國家：{d.get('countryCode','')}\nISP：{d.get('isp','')}"
            return f"查詢失敗：{resp.status_code}"
        except Exception as e:
            return f"查詢失敗：{e}"

    return f"未知動作：{action}"


def execute_auto_skill(action, goal="", skill_name="", code="", test_input=""):
    """自動技能生成與部署：用AI寫新技能並部署到bot"""

    if action == "generate":
        try:
            from anthropic import Anthropic
            c = Anthropic()
            resp = c.messages.create(
                model="claude-sonnet-4-6",
                max_tokens=2000,
                messages=[{"role": "user", "content": f"""你是Python專家。請為小牛馬Telegram Bot生成一個新的execute_函數。

技能需求：{goal}
函數名稱：execute_{skill_name or 'new_skill'}

要求：
1. 完整錯誤處理
2. 回傳值必須是字串
3. 使用已安裝的標準庫（requests, json, os, pathlib等）
4. 只輸出Python程式碼，不要任何說明

輸出完整函數程式碼："""}]
            )
            generated = resp.content[0].text
            draft = Path(__file__).parent / f"skill_{skill_name or 'draft'}.py"
            draft.write_text(generated, encoding="utf-8")
            return f"✅ 技能已生成\n儲存至：{draft}\n\n```python\n{generated[:600]}\n```"
        except Exception as e:
            return f"生成失敗：{e}"

    elif action == "test":
        try:
            compile(code, "<string>", "exec")
            return "✅ 程式碼語法正確"
        except SyntaxError as e:
            return f"❌ 語法錯誤：{e}"

    elif action == "deploy":
        try:
            if not code:
                draft = Path(__file__).parent / f"skill_{skill_name}.py"
                if draft.exists():
                    code = draft.read_text(encoding="utf-8")
                else:
                    return "請提供程式碼"
            bot_path = Path(__file__)
            content = bot_path.read_text(encoding="utf-8")
            marker = "\nasync def start("
            if marker in content:
                new_content = content.replace(marker, f"\n\n# ── 自動部署：{skill_name} ──\n{code}\n{marker}")
                bot_path.write_text(new_content, encoding="utf-8")
                return f"✅ 技能 {skill_name} 已部署，重啟後生效"
            return "部署失敗：找不到插入點"
        except Exception as e:
            return f"部署失敗：{e}"

    elif action == "list_skills":
        import re
        content = Path(__file__).read_text(encoding="utf-8")
        skills = re.findall(r'def execute_(\w+)\(', content)
        return f"已部署技能：{len(skills)} 個\n" + "\n".join(f"• {s}" for s in sorted(set(skills)))

    return f"未知動作：{action}"


def execute_smart_home(action, device="", value="", host="", token=""):
    """智慧家居控制：Home Assistant API整合"""
    ha_host = host or os.environ.get("HA_HOST", "http://homeassistant.local:8123")
    ha_token = token or os.environ.get("HA_TOKEN", "")
    headers = {"Authorization": f"Bearer {ha_token}", "Content-Type": "application/json"}

    if action == "list_devices":
        try:
            resp = requests.get(f"{ha_host}/api/states", headers=headers, timeout=10)
            if resp.status_code == 200:
                states = resp.json()
                result = [f"🏠 智慧家居設備（{len(states)} 個）："]
                for s in states[:25]:
                    result.append(f"  {s['entity_id']}: {s['state']}")
                return "\n".join(result)
            return f"連線失敗 {resp.status_code} — 請設定 HA_HOST 和 HA_TOKEN"
        except Exception as e:
            return f"連線失敗：{e}\n請設定環境變數 HA_HOST 和 HA_TOKEN"

    elif action in ("turn_on", "turn_off"):
        try:
            domain = device.split(".")[0] if "." in device else "homeassistant"
            svc = "turn_on" if action == "turn_on" else "turn_off"
            resp = requests.post(f"{ha_host}/api/services/{domain}/{svc}",
                                 headers=headers, json={"entity_id": device}, timeout=10)
            return f"✅ {device} {'開啟' if action=='turn_on' else '關閉'}" if resp.status_code == 200 else f"失敗：{resp.status_code}"
        except Exception as e:
            return f"控制失敗：{e}"

    elif action == "get_state":
        try:
            resp = requests.get(f"{ha_host}/api/states/{device}", headers=headers, timeout=10)
            if resp.status_code == 200:
                s = resp.json()
                return f"{device}\n狀態：{s['state']}\n屬性：{json.dumps(s.get('attributes',{}), ensure_ascii=False)[:300]}"
            return f"查詢失敗：{resp.status_code}"
        except Exception as e:
            return f"查詢失敗：{e}"

    elif action == "set_value":
        try:
            domain = device.split(".")[0] if "." in device else "input_number"
            resp = requests.post(f"{ha_host}/api/services/{domain}/set_value",
                                 headers=headers, json={"entity_id": device, "value": value}, timeout=10)
            return f"✅ {device} 設定為 {value}" if resp.status_code == 200 else f"失敗：{resp.status_code}"
        except Exception as e:
            return f"設定失敗：{e}"

    elif action == "run_scene":
        try:
            resp = requests.post(f"{ha_host}/api/services/scene/turn_on",
                                 headers=headers, json={"entity_id": f"scene.{device}"}, timeout=10)
            return f"✅ 場景 {device} 已啟動" if resp.status_code == 200 else f"失敗：{resp.status_code}"
        except Exception as e:
            return f"執行失敗：{e}"

    return f"未知動作：{action}"


def execute_goal_manager(action, goal="", goal_id="", steps="", priority="normal"):
    """目標管理系統：設定長期目標、AI拆解步驟、追蹤執行進度"""
    GOALS_DB = Path(__file__).parent / "goals.db"

    def init_db():
        conn = sqlite3.connect(GOALS_DB)
        conn.execute("""CREATE TABLE IF NOT EXISTS goals (
            id INTEGER PRIMARY KEY AUTOINCREMENT,
            title TEXT NOT NULL,
            steps TEXT,
            status TEXT DEFAULT 'pending',
            priority TEXT DEFAULT 'normal',
            progress INTEGER DEFAULT 0,
            created_at TIMESTAMP DEFAULT CURRENT_TIMESTAMP,
            updated_at TIMESTAMP DEFAULT CURRENT_TIMESTAMP
        )""")
        conn.commit()
        conn.close()

    init_db()

    if action == "add":
        try:
            if not steps:
                from anthropic import Anthropic
                c = Anthropic()
                resp = c.messages.create(
                    model="claude-sonnet-4-6", max_tokens=800,
                    messages=[{"role": "user", "content": f"將以下目標分解為5-8個具體執行步驟，每步驟一行，格式：1. 步驟\n\n目標：{goal}"}]
                )
                steps = resp.content[0].text
            conn = sqlite3.connect(GOALS_DB)
            cur = conn.execute("INSERT INTO goals (title, steps, priority) VALUES (?, ?, ?)", (goal, steps, priority))
            gid = cur.lastrowid
            conn.commit()
            conn.close()
            return f"✅ 目標建立 (ID:{gid})\n{goal}\n\n步驟：\n{steps}"
        except Exception as e:
            return f"建立失敗：{e}"

    elif action == "list":
        conn = sqlite3.connect(GOALS_DB)
        rows = conn.execute("SELECT id, title, status, priority, progress FROM goals ORDER BY id DESC").fetchall()
        conn.close()
        if not rows:
            return "目前沒有任何目標"
        icons = {"pending": "⏳", "in_progress": "🔄", "completed": "✅", "failed": "❌"}
        return "📋 目標清單：\n" + "\n".join(f"  {icons.get(r[2],'❓')} [{r[0]}] {r[1]} ({r[3]}) {r[4]}%" for r in rows)

    elif action == "detail":
        conn = sqlite3.connect(GOALS_DB)
        row = conn.execute("SELECT * FROM goals WHERE id=?", (goal_id,)).fetchone()
        conn.close()
        if not row:
            return f"找不到目標 {goal_id}"
        return f"目標[{row[0]}]：{row[1]}\n狀態：{row[3]} | 優先：{row[4]} | 進度：{row[5]}%\n\n步驟：\n{row[2]}"

    elif action == "update_status":
        conn = sqlite3.connect(GOALS_DB)
        conn.execute("UPDATE goals SET status=?, updated_at=CURRENT_TIMESTAMP WHERE id=?", (steps, goal_id))
        conn.commit()
        conn.close()
        return f"✅ 目標 {goal_id} 狀態 → {steps}"

    elif action == "set_progress":
        conn = sqlite3.connect(GOALS_DB)
        conn.execute("UPDATE goals SET progress=?, updated_at=CURRENT_TIMESTAMP WHERE id=?", (int(steps or 0), goal_id))
        conn.commit()
        conn.close()
        return f"✅ 目標 {goal_id} 進度 → {steps}%"

    elif action == "delete":
        conn = sqlite3.connect(GOALS_DB)
        conn.execute("DELETE FROM goals WHERE id=?", (goal_id,))
        conn.commit()
        conn.close()
        return f"✅ 目標 {goal_id} 已刪除"

    elif action == "next_task":
        conn = sqlite3.connect(GOALS_DB)
        row = conn.execute("SELECT * FROM goals WHERE status IN ('pending','in_progress') ORDER BY priority DESC, id ASC LIMIT 1").fetchone()
        conn.close()
        if not row:
            return "所有目標已完成或無目標"
        return f"🤖 下一個待執行目標：\n[{row[0]}] {row[1]}\n\n步驟：\n{row[2][:500] if row[2] else '待規劃'}"

    return f"未知動作：{action}"


def execute_auto_trade(action, symbol="", amount=0.0, price=0.0, order_type="market", api_key="", api_secret=""):
    """加密貨幣自動交易：Binance API"""
    bk = api_key or os.environ.get("BINANCE_KEY", "")
    bs = api_secret or os.environ.get("BINANCE_SECRET", "")

    if action == "price":
        try:
            sym = symbol.upper().replace("/", "").replace("-", "")
            resp = requests.get(f"https://api.binance.com/api/v3/ticker/24hr?symbol={sym}", timeout=5)
            if resp.status_code == 200:
                d = resp.json()
                chg = float(d.get("priceChangePercent", "0"))
                arrow = "▲" if chg >= 0 else "▼"
                return f"{sym} 行情：\n價格：{d.get('lastPrice','')} USDT\n24h：{arrow} {chg}%\n高：{d.get('highPrice','')} | 低：{d.get('lowPrice','')} | 量：{d.get('volume','')}"
            return f"查詢失敗：{resp.status_code}"
        except Exception as e:
            return f"查詢失敗：{e}"

    elif action == "balance":
        if not bk or not bs:
            return "需要設定 BINANCE_KEY 和 BINANCE_SECRET"
        try:
            import hmac, hashlib
            ts = int(time.time() * 1000)
            params = f"timestamp={ts}"
            sig = hmac.new(bs.encode(), params.encode(), hashlib.sha256).hexdigest()
            resp = requests.get(f"https://api.binance.com/api/v3/account?{params}&signature={sig}",
                                headers={"X-MBX-APIKEY": bk}, timeout=10)
            if resp.status_code == 200:
                non_zero = [b for b in resp.json().get("balances", []) if float(b["free"]) > 0.0001]
                return "💰 帳戶餘額：\n" + "\n".join(f"  {b['asset']}: {b['free']}" for b in non_zero[:15])
            return f"查詢失敗：{resp.text}"
        except Exception as e:
            return f"查詢失敗：{e}"

    elif action in ("buy", "sell"):
        if not bk or not bs:
            return "需要設定 BINANCE_KEY 和 BINANCE_SECRET 才能下單"
        try:
            import hmac, hashlib
            sym = symbol.upper().replace("/", "").replace("-", "")
            side = "BUY" if action == "buy" else "SELL"
            ts = int(time.time() * 1000)
            params = f"symbol={sym}&side={side}&type={order_type.upper()}&quantity={amount}&timestamp={ts}"
            if order_type.lower() == "limit" and float(price) > 0:
                params += f"&price={price}&timeInForce=GTC"
            sig = hmac.new(bs.encode(), params.encode(), hashlib.sha256).hexdigest()
            resp = requests.post(f"https://api.binance.com/api/v3/order?{params}&signature={sig}",
                                 headers={"X-MBX-APIKEY": bk}, timeout=10)
            if resp.status_code == 200:
                d = resp.json()
                return f"✅ 訂單成功\nID: {d.get('orderId','')}\n{side} {amount} {sym} @ {order_type}\n狀態：{d.get('status','')}"
            return f"下單失敗：{resp.text}"
        except Exception as e:
            return f"下單失敗：{e}"

    elif action == "open_orders":
        if not bk or not bs:
            return "需要設定 BINANCE_KEY 和 BINANCE_SECRET"
        try:
            import hmac, hashlib
            ts = int(time.time() * 1000)
            sym_part = f"&symbol={symbol.upper()}" if symbol else ""
            params = f"timestamp={ts}{sym_part}"
            sig = hmac.new(bs.encode(), params.encode(), hashlib.sha256).hexdigest()
            resp = requests.get(f"https://api.binance.com/api/v3/openOrders?{params}&signature={sig}",
                                headers={"X-MBX-APIKEY": bk}, timeout=10)
            if resp.status_code == 200:
                orders = resp.json()
                if not orders:
                    return "目前無掛單"
                return f"掛單列表（{len(orders)}筆）：\n" + "\n".join(
                    f"  [{o['orderId']}] {o['side']} {o['origQty']} {o['symbol']} @ {o['price']}" for o in orders[:10])
            return f"查詢失敗：{resp.text}"
        except Exception as e:
            return f"查詢失敗：{e}"

    return f"未知動作：{action}"


def execute_knowledge_base(action, content="", query="", tag="", kb_id=""):
    """知識庫：儲存、全文搜尋、管理結構化知識"""
    KB_DB = Path(__file__).parent / "knowledge_base.db"

    def init_kb():
        conn = sqlite3.connect(KB_DB)
        conn.execute("""CREATE TABLE IF NOT EXISTS knowledge (
            id INTEGER PRIMARY KEY AUTOINCREMENT,
            content TEXT NOT NULL,
            summary TEXT,
            tags TEXT DEFAULT '',
            created_at TIMESTAMP DEFAULT CURRENT_TIMESTAMP
        )""")
        conn.commit()
        conn.close()

    init_kb()

    if action == "add":
        try:
            summary = content[:150] if len(content) > 150 else content
            conn = sqlite3.connect(KB_DB)
            cur = conn.execute("INSERT INTO knowledge (content, summary, tags) VALUES (?, ?, ?)",
                               (content, summary, tag))
            kid = cur.lastrowid
            conn.commit()
            conn.close()
            return f"✅ 知識儲存 (ID:{kid})\n標籤：{tag}\n摘要：{summary[:100]}"
        except Exception as e:
            return f"儲存失敗：{e}"

    elif action == "search":
        try:
            conn = sqlite3.connect(KB_DB)
            rows = conn.execute(
                "SELECT id, summary, tags, created_at FROM knowledge WHERE content LIKE ? OR tags LIKE ? ORDER BY id DESC LIMIT 10",
                (f"%{query}%", f"%{query}%")
            ).fetchall()
            conn.close()
            if not rows:
                return f"找不到「{query}」相關知識"
            return f"🔍 找到 {len(rows)} 筆：\n" + "\n".join(f"  [{r[0]}] {r[1][:100]} [{r[2]}]" for r in rows)
        except Exception as e:
            return f"搜尋失敗：{e}"

    elif action == "get":
        conn = sqlite3.connect(KB_DB)
        row = conn.execute("SELECT * FROM knowledge WHERE id=?", (kb_id,)).fetchone()
        conn.close()
        if not row:
            return f"找不到知識 {kb_id}"
        return f"[{row[0]}] {row[3] or '無標籤'} | {row[4]}\n\n{row[1]}"

    elif action == "list":
        conn = sqlite3.connect(KB_DB)
        rows = conn.execute("SELECT id, summary, tags, created_at FROM knowledge ORDER BY id DESC LIMIT 20").fetchall()
        total = conn.execute("SELECT COUNT(*) FROM knowledge").fetchone()[0]
        conn.close()
        if not rows:
            return "知識庫空白"
        return f"📚 知識庫（共{total}條）：\n" + "\n".join(f"  [{r[0]}] {r[1][:80]} [{r[2]}]" for r in rows)

    elif action == "delete":
        conn = sqlite3.connect(KB_DB)
        conn.execute("DELETE FROM knowledge WHERE id=?", (kb_id,))
        conn.commit()
        conn.close()
        return f"✅ 知識 {kb_id} 已刪除"

    elif action == "stats":
        conn = sqlite3.connect(KB_DB)
        total = conn.execute("SELECT COUNT(*) FROM knowledge").fetchone()[0]
        tags = conn.execute("SELECT tags, COUNT(*) FROM knowledge GROUP BY tags ORDER BY COUNT(*) DESC LIMIT 10").fetchall()
        conn.close()
        return f"📊 知識庫統計\n總條目：{total}\n\n標籤分佈：\n" + "\n".join(f"  {t[0] or '無標籤'}: {t[1]}條" for t in tags)

    return f"未知動作：{action}"


def execute_emotion_detect(action, text="", image_path=""):
    """情緒偵測：從文字/臉部偵測情緒狀態"""

    if action == "from_text":
        try:
            pos_words = ["開心","高興","棒","好","讚","喜歡","愛","感謝","謝謝","哈哈","😊","😄","👍","❤️","爽","讚讚"]
            neg_words = ["難過","生氣","討厭","煩","累","痛","哭","傷心","憤怒","失望","😢","😡","😤","💔","爛","幹"]
            pos = sum(1 for w in pos_words if w in text)
            neg = sum(1 for w in neg_words if w in text)
            basic = "正面" if pos > neg else ("負面" if neg > pos else "中性")

            from anthropic import Anthropic
            c = Anthropic()
            resp = c.messages.create(
                model="claude-haiku-4-5-20251001", max_tokens=200,
                messages=[{"role": "user", "content": f"分析情緒（一行一項）：主要情緒、強度1-10、建議回應\n文字：{text}"}]
            )
            return f"情緒分析\n規則：{basic}（正面{pos}負面{neg}）\n\nAI分析：\n{resp.content[0].text}"
        except Exception as e:
            return f"分析失敗：{e}"

    elif action == "from_face":
        try:
            import base64, io as _io
            from PIL import Image as PILImg
            if image_path:
                img = PILImg.open(image_path)
            else:
                import pyautogui
                img = pyautogui.screenshot()
            buf = _io.BytesIO()
            img.save(buf, format="JPEG")
            b64 = base64.b64encode(buf.getvalue()).decode()

            from anthropic import Anthropic
            c = Anthropic()
            resp = c.messages.create(
                model="claude-sonnet-4-6", max_tokens=300,
                messages=[{"role": "user", "content": [
                    {"type": "image", "source": {"type": "base64", "media_type": "image/jpeg", "data": b64}},
                    {"type": "text", "text": "分析圖片中人臉情緒：臉數、每張臉的情緒（開心/悲傷/憤怒/驚訝/中性）、整體氛圍"}
                ]}]
            )
            return f"臉部情緒分析：\n{resp.content[0].text}"
        except Exception as e:
            return f"臉部分析失敗：{e}"

    return f"未知動作：{action}"


def execute_voice_id(action, name="", audio_path="", duration=5):
    """聲紋辨識：登記/辨識說話者身份"""
    VOICE_DIR = Path(__file__).parent / "voice_profiles"
    VOICE_DIR.mkdir(exist_ok=True)
    META_FILE = Path(__file__).parent / "voice_profiles.json"

    def load_meta():
        return json.loads(META_FILE.read_text("utf-8")) if META_FILE.exists() else {}

    def save_meta(d):
        META_FILE.write_text(json.dumps(d, ensure_ascii=False, indent=2), "utf-8")

    if action == "enroll":
        try:
            import sounddevice as sd
            import scipy.io.wavfile as wf
            import numpy as np
            sr = 16000
            rec = sd.rec(int(float(duration) * sr), samplerate=sr, channels=1)
            sd.wait()
            p = VOICE_DIR / f"{name}.wav"
            wf.write(str(p), sr, rec)
            meta = load_meta()
            meta[name] = str(p)
            save_meta(meta)
            return f"✅ {name} 聲紋已登記"
        except ImportError:
            return "需要安裝：pip install sounddevice scipy"
        except Exception as e:
            return f"登記失敗：{e}"

    elif action == "identify":
        try:
            import librosa
            import numpy as np
            meta = load_meta()
            if not meta:
                return "尚未登記任何聲紋"
            if not audio_path:
                return "請提供音訊檔案路徑"
            q_audio, _ = librosa.load(audio_path, sr=16000)
            q_mfcc = librosa.feature.mfcc(y=q_audio, sr=16000, n_mfcc=13).mean(axis=1)
            best, best_score = None, float('inf')
            for person, path in meta.items():
                p_audio, _ = librosa.load(path, sr=16000)
                p_mfcc = librosa.feature.mfcc(y=p_audio, sr=16000, n_mfcc=13).mean(axis=1)
                score = np.linalg.norm(q_mfcc - p_mfcc)
                if score < best_score:
                    best_score, best = score, person
            conf = max(0, 100 - int(best_score * 10))
            return f"聲紋辨識：{best}（信心度 {conf}%）"
        except ImportError:
            return "需要安裝：pip install librosa"
        except Exception as e:
            return f"辨識失敗：{e}"

    elif action == "list":
        meta = load_meta()
        return "已登記聲紋：\n" + "\n".join(f"• {n}" for n in meta) if meta else "尚未登記任何聲紋"

    elif action == "delete":
        meta = load_meta()
        if name in meta:
            Path(meta[name]).unlink(missing_ok=True)
            del meta[name]
            save_meta(meta)
            return f"✅ {name} 聲紋已刪除"
        return f"找不到 {name}"

    return f"未知動作：{action}"


def execute_pentest(action, target="", port_range="1-1000", timeout=2):
    """滲透測試：針對自己網路的安全評估工具"""
    import socket, threading

    if action == "port_scan":
        open_ports = []
        try:
            target_ip = socket.gethostbyname(target)
        except Exception:
            target_ip = target
        if "-" in str(port_range):
            start, end = map(int, str(port_range).split("-"))
            ports = list(range(start, min(end + 1, start + 500)))
        else:
            ports = [int(p) for p in str(port_range).split(",")]

        def scan(port):
            try:
                s = socket.socket()
                s.settimeout(float(timeout))
                if s.connect_ex((target_ip, port)) == 0:
                    try:
                        svc = socket.getservbyport(port)
                    except Exception:
                        svc = "unknown"
                    open_ports.append((port, svc))
                s.close()
            except Exception:
                pass

        threads = [threading.Thread(target=scan, args=(p,)) for p in ports]
        for t in threads:
            t.start()
        for t in threads:
            t.join(timeout=1)
        open_ports.sort()
        return f"🔍 埠掃描 {target_ip}\n開放：{len(open_ports)} 個\n" + "\n".join(f"  {p}/tcp {s}" for p, s in open_ports[:20])

    elif action == "ssl_check":
        import ssl
        try:
            ctx = ssl.create_default_context()
            with ctx.wrap_socket(socket.socket(), server_hostname=target) as s:
                s.settimeout(5)
                s.connect((target, 443))
                cert = s.getpeercert()
            import datetime
            exp = datetime.datetime.strptime(cert['notAfter'], '%b %d %H:%M:%S %Y %Z')
            days = (exp - datetime.datetime.utcnow()).days
            return f"SSL憑證：{target}\n到期：{cert['notAfter']}\n剩餘：{days}天\n{'⚠️ 即將到期！' if days < 30 else '✅ 有效'}"
        except Exception as e:
            return f"SSL檢查失敗：{e}"

    elif action == "http_headers":
        try:
            url = f"https://{target}" if not target.startswith("http") else target
            resp = requests.get(url, timeout=5, verify=False, allow_redirects=False)
            sec_headers = ["Strict-Transport-Security", "X-Content-Type-Options", "X-Frame-Options",
                           "Content-Security-Policy", "X-XSS-Protection", "Referrer-Policy"]
            result = [f"🔒 安全標頭：{target}"]
            for h in sec_headers:
                if h in resp.headers:
                    result.append(f"  ✅ {h}")
                else:
                    result.append(f"  ❌ 缺少 {h}")
            return "\n".join(result)
        except Exception as e:
            return f"分析失敗：{e}"

    elif action == "vuln_scan":
        try:
            base = f"https://{target}" if not target.startswith("http") else target
            paths = [".env", ".git/config", "wp-config.php", "phpinfo.php", "admin/", "backup.sql", ".DS_Store"]
            found = []
            for path in paths:
                try:
                    r = requests.get(f"{base}/{path}", timeout=3, verify=False, allow_redirects=False)
                    if r.status_code == 200:
                        found.append(f"  ⚠️ /{path} 可存取！")
                except Exception:
                    pass
            if found:
                return f"🚨 發現 {len(found)} 個潛在漏洞：\n" + "\n".join(found)
            return f"✅ {target} 未發現常見敏感路徑暴露"
        except Exception as e:
            return f"掃描失敗：{e}"

    elif action == "password_audit":
        import re
        results = []
        for pwd in target.split(","):
            pwd = pwd.strip()
            score = 0
            issues = []
            if len(pwd) >= 12: score += 2
            elif len(pwd) >= 8: score += 1
            else: issues.append("太短")
            if re.search(r'[A-Z]', pwd): score += 1
            else: issues.append("缺大寫")
            if re.search(r'[a-z]', pwd): score += 1
            else: issues.append("缺小寫")
            if re.search(r'\d', pwd): score += 1
            else: issues.append("缺數字")
            if re.search(r'[!@#$%^&*]', pwd): score += 2
            else: issues.append("缺符號")
            strength = ["極弱","弱","普通","強","極強"][min(score//2, 4)]
            results.append(f"'{pwd[:3]}***': {strength}({score}/7) {' '.join(issues)}")
        return "密碼強度審計：\n" + "\n".join(results)

    return f"未知動作：{action}"


def execute_proactive_alert(action, name="", condition="", threshold="", target="",
                             interval=60, chat_id=None, _bot_send=None):
    """主動預警系統：持續監控條件並自動通知"""
    import threading
    ALERTS_FILE = Path(__file__).parent / "proactive_alerts.json"

    def load_alerts():
        return json.loads(ALERTS_FILE.read_text("utf-8")) if ALERTS_FILE.exists() else {}

    def save_alerts(d):
        ALERTS_FILE.write_text(json.dumps(d, ensure_ascii=False, indent=2), "utf-8")

    if action == "add":
        alerts = load_alerts()
        alerts[name] = {"condition": condition, "threshold": threshold, "target": target,
                        "interval": interval, "active": True, "triggered": 0}
        save_alerts(alerts)
        return f"✅ 預警 '{name}' 建立\n條件：{condition} {threshold}\n目標：{target}\n間隔：{interval}秒"

    elif action == "list":
        alerts = load_alerts()
        if not alerts:
            return "無預警設定"
        return "🚨 預警清單：\n" + "\n".join(
            f"  {'🟢' if a.get('active') else '🔴'} [{n}] {a['condition']} {a.get('threshold','')} (觸發{a.get('triggered',0)}次)"
            for n, a in alerts.items())

    elif action == "delete":
        alerts = load_alerts()
        if name in alerts:
            del alerts[name]
            save_alerts(alerts)
            return f"✅ 預警 '{name}' 已刪除"
        return f"找不到 '{name}'"

    elif action == "toggle":
        alerts = load_alerts()
        if name in alerts:
            alerts[name]["active"] = not alerts[name].get("active", True)
            save_alerts(alerts)
            return f"✅ 預警 '{name}' {'啟用' if alerts[name]['active'] else '停用'}"
        return f"找不到 '{name}'"

    elif action == "start_all":
        alerts = load_alerts()

        def monitor_loop():
            while True:
                current = load_alerts()
                for aname, acfg in current.items():
                    if not acfg.get("active"):
                        continue
                    try:
                        msg = None
                        cond = acfg["condition"]
                        thresh = acfg.get("threshold", "")
                        tgt = acfg.get("target", "")

                        if cond == "cpu_above":
                            import psutil
                            val = psutil.cpu_percent(interval=1)
                            if val > float(thresh):
                                msg = f"🚨 CPU {val}% > {thresh}%"
                        elif cond == "memory_above":
                            import psutil
                            val = psutil.virtual_memory().percent
                            if val > float(thresh):
                                msg = f"🚨 記憶體 {val}% > {thresh}%"
                        elif cond == "price_above":
                            r = requests.get(f"https://api.binance.com/api/v3/ticker/price?symbol={tgt}", timeout=5)
                            val = float(r.json().get("price", 0))
                            if val > float(thresh):
                                msg = f"📈 {tgt} 價格 {val} > {thresh}"
                        elif cond == "price_below":
                            r = requests.get(f"https://api.binance.com/api/v3/ticker/price?symbol={tgt}", timeout=5)
                            val = float(r.json().get("price", 0))
                            if val < float(thresh):
                                msg = f"📉 {tgt} 價格 {val} < {thresh}"
                        elif cond == "news_keyword":
                            import feedparser
                            feed = feedparser.parse(f"https://news.google.com/rss/search?q={urllib.parse.quote(tgt)}&hl=zh-TW")
                            if feed.entries:
                                msg = f"📰 新聞 [{tgt}]\n{feed.entries[0].get('title','')}"

                        if msg and _bot_send and chat_id:
                            import asyncio
                            asyncio.run_coroutine_threadsafe(
                                _bot_send(chat_id=chat_id, text=f"[預警:{aname}] {msg}"),
                                asyncio.get_event_loop()
                            )
                            cur_alerts = load_alerts()
                            if aname in cur_alerts:
                                cur_alerts[aname]["triggered"] = cur_alerts[aname].get("triggered", 0) + 1
                                save_alerts(cur_alerts)
                    except Exception:
                        pass
                time.sleep(float(interval) if interval else 60)

        threading.Thread(target=monitor_loop, daemon=True).start()
        return f"✅ 已啟動 {len(alerts)} 個預警監控"

    return f"未知動作：{action}"


def execute_multi_deploy(action, remote_host="", remote_user="", remote_pass="", remote_path="/tmp/niu_bot"):
    """多機器部署：將bot部署到遠端伺服器同步運行"""
    try:
        import paramiko
    except ImportError:
        return "需要安裝 paramiko: pip install paramiko"

    def get_ssh():
        ssh = paramiko.SSHClient()
        ssh.set_missing_host_key_policy(paramiko.AutoAddPolicy())
        ssh.connect(remote_host, username=remote_user, password=remote_pass, timeout=15)
        return ssh

    if action == "deploy":
        try:
            ssh = get_ssh()
            sftp = ssh.open_sftp()
            bot_path = Path(__file__)
            env_path = bot_path.parent / ".env"

            # 建立目錄
            ssh.exec_command(f"mkdir -p {remote_path}")
            time.sleep(0.5)
            sftp.put(str(bot_path), f"{remote_path}/bot.py")
            if env_path.exists():
                sftp.put(str(env_path), f"{remote_path}/.env")
            sftp.close()

            for cmd in [
                "pip install python-telegram-bot anthropic requests python-dotenv feedparser -q",
                f"pkill -f '{remote_path}/bot.py' 2>/dev/null; true",
                f"nohup python {remote_path}/bot.py > {remote_path}/bot.log 2>&1 &"
            ]:
                ssh.exec_command(cmd)
                time.sleep(1)
            ssh.close()
            return f"✅ Bot 已部署到 {remote_host}:{remote_path}"
        except Exception as e:
            return f"部署失敗：{e}"

    elif action == "status":
        try:
            ssh = get_ssh()
            _, stdout, _ = ssh.exec_command(f"pgrep -f '{remote_path}/bot.py' && echo RUNNING || echo STOPPED")
            status = stdout.read().decode().strip()
            ssh.close()
            return f"{remote_host} Bot：{status}"
        except Exception as e:
            return f"查詢失敗：{e}"

    elif action == "sync":
        try:
            ssh = get_ssh()
            sftp = ssh.open_sftp()
            sftp.put(str(Path(__file__)), f"{remote_path}/bot.py")
            sftp.close()
            ssh.exec_command(f"pkill -f '{remote_path}/bot.py'; sleep 1; nohup python {remote_path}/bot.py > {remote_path}/bot.log 2>&1 &")
            ssh.close()
            return f"✅ 技能已同步到 {remote_host} 並重啟"
        except Exception as e:
            return f"同步失敗：{e}"

    elif action == "log":
        try:
            ssh = get_ssh()
            _, stdout, _ = ssh.exec_command(f"tail -50 {remote_path}/bot.log")
            log = stdout.read().decode()
            ssh.close()
            return f"遠端日誌 ({remote_host}):\n{log[-2000:]}"
        except Exception as e:
            return f"讀取日誌失敗：{e}"

    return f"未知動作：{action}"


def execute_self_benchmark(action):
    """自我評估：測試各功能健康狀態並產生報告"""
    if action == "run":
        tests = {
            "網路連線": lambda: requests.get("https://google.com", timeout=3).status_code == 200,
            "Claude API Key": lambda: bool(os.environ.get("ANTHROPIC_API_KEY", "")),
            "Telegram Token": lambda: bool(os.environ.get("BOT_TOKEN", "")),
            "資料庫": lambda: sqlite3.connect(Path(__file__).parent / "memory.db").execute("SELECT 1").fetchone() is not None,
            "截圖功能": lambda: bool(__import__("pyautogui").screenshot()),
            "HuggingFace Key": lambda: bool(os.environ.get("HF_TOKEN", "")),
            "知識庫": lambda: (Path(__file__).parent / "knowledge_base.db").exists(),
            "目標管理": lambda: (Path(__file__).parent / "goals.db").exists(),
            "Binance Key": lambda: bool(os.environ.get("BINANCE_KEY", "")),
            "HA Token": lambda: bool(os.environ.get("HA_TOKEN", "")),
            "VirusTotal Key": lambda: bool(os.environ.get("VIRUSTOTAL_KEY", "")),
            "預警設定": lambda: (Path(__file__).parent / "proactive_alerts.json").exists(),
        }
        passed = 0
        lines = ["🤖 自我評估報告\n"]
        for name, test in tests.items():
            try:
                ok = test()
                lines.append(f"  {'✅' if ok else '⚠️'} {name}")
                if ok:
                    passed += 1
            except Exception as e:
                lines.append(f"  ❌ {name}: {str(e)[:30]}")
        score = int(passed / len(tests) * 100)
        lines.insert(1, f"健康度：{score}% ({passed}/{len(tests)})")
        return "\n".join(lines)

    elif action == "skill_count":
        import re
        content = Path(__file__).read_text(encoding="utf-8")
        skills = sorted(set(re.findall(r'def execute_(\w+)\(', content)))
        return f"技能函數總數：{len(skills)}\n\n" + "\n".join(f"• {s}" for s in skills)

    elif action == "memory_stats":
        db = Path(__file__).parent / "memory.db"
        if not db.exists():
            return "資料庫不存在"
        conn = sqlite3.connect(db)
        hist = conn.execute("SELECT COUNT(*) FROM chat_history").fetchone()[0]
        memo = conn.execute("SELECT COUNT(*) FROM long_term_memory").fetchone()[0]
        conn.close()
        kb_db = Path(__file__).parent / "knowledge_base.db"
        kb_count = 0
        if kb_db.exists():
            c2 = sqlite3.connect(kb_db)
            kb_count = c2.execute("SELECT COUNT(*) FROM knowledge").fetchone()[0]
            c2.close()
        return f"記憶統計：\n對話歷史：{hist} 條\n長期記憶：{memo} 條\n知識庫：{kb_count} 條"

    return f"未知動作：{action}"


async def start(update: Update, context: ContextTypes.DEFAULT_TYPE):
    await update.message.reply_text("你好！我是小牛馬，有什麼可以幫你的？（我還記得我們之前聊過的事 😄）")


async def handle_voice_message(update: Update, context: ContextTypes.DEFAULT_TYPE):
    """處理 Telegram 語音訊息：OGG → WAV → STT → Claude → 語音回覆"""
    try:
        chat_id = update.effective_chat.id
        is_owner = update.effective_user.id == OWNER_ID
        sender_name = "于晏" if is_owner else (update.effective_user.first_name or str(update.effective_user.id))

        await context.bot.send_chat_action(chat_id=chat_id, action="typing")

        # 下載語音檔（OGG OPUS）
        voice = update.message.voice or update.message.audio
        if not voice:
            return
        voice_file = await context.bot.get_file(voice.file_id)
        import tempfile, os
        with tempfile.NamedTemporaryFile(suffix=".ogg", delete=False) as tmp_ogg:
            ogg_path = tmp_ogg.name
        await voice_file.download_to_drive(ogg_path)

        # OGG → WAV（用 ffmpeg）
        import imageio_ffmpeg
        ffmpeg = imageio_ffmpeg.get_ffmpeg_exe()
        wav_path = ogg_path.replace(".ogg", ".wav")
        import subprocess
        subprocess.run([ffmpeg, "-y", "-i", ogg_path, "-ar", "16000", "-ac", "1", wav_path],
                       capture_output=True)
        Path(ogg_path).unlink(missing_ok=True)

        # WAV → 文字（Google STT）
        import speech_recognition as sr
        recognizer = sr.Recognizer()
        user_text = ""
        try:
            with sr.AudioFile(wav_path) as source:
                audio = recognizer.record(source)
            user_text = recognizer.recognize_google(audio, language="zh-TW")
        except sr.UnknownValueError:
            user_text = ""
        except Exception as e:
            user_text = ""
        Path(wav_path).unlink(missing_ok=True)

        if not user_text:
            await update.message.reply_text("🎤 聽不清楚，請再說一次～")
            return

        # 顯示辨識結果
        await update.message.reply_text(f"🎤 你說：{user_text}")
        log_message("🎤", sender_name, chat_id, user_text)

        # 傳給 Claude 處理（同 handle_message 邏輯）
        saved_text = f"[{sender_name}]: {user_text}" if update.effective_chat.type in ("group","supergroup") else user_text
        save_message(chat_id, "user", saved_text)
        history = load_history(chat_id)[-40:]

        base_system = SYSTEM_PROMPT_OWNER if is_owner else SYSTEM_PROMPT_DEFAULT
        _voice_is_group = update.effective_chat.type in ("group", "supergroup")
        if _voice_is_group:
            if is_owner:
                base_system += f"\n\n【當前說話者】{sender_name}（主人，回覆結尾請叫于晏哥）"
            else:
                base_system += f"\n\n【當前說話者】{sender_name}（不是主人，用「{sender_name}」稱呼，絕對不要叫他于晏哥）"
        memories = load_long_term_memory(chat_id)
        if memories:
            mem_text = "\n".join(f"- [{m['id']}] {m['content']}" for m in memories)
            system = base_system + f"\n\n【長期記憶】\n{mem_text}"
        else:
            system = base_system

        response = client.messages.create(
            model="claude-sonnet-4-6",
            max_tokens=1024,
            system=system,
            messages=history
        )

        reply_text = response.content[0].text if response.content else "..."
        if not is_owner:
            reply_text = _fix_group_reply(reply_text, sender_name)
        save_message(chat_id, "assistant", reply_text)
        log_message("<<", "小牛馬", chat_id, reply_text)

        # 用語音回覆
        await context.bot.send_chat_action(chat_id=chat_id, action="record_voice")
        import asyncio, io as _io
        loop = asyncio.get_running_loop()
        try:
            ogg_data = await loop.run_in_executor(None, generate_voice_ogg, reply_text, "zh-TW-YunJheNeural")
            await update.message.reply_voice(voice=_io.BytesIO(ogg_data))
        except Exception:
            await update.message.reply_text(reply_text)

    except Exception as e:
        await update.message.reply_text(f"❌ 語音處理失敗：{e}")


def _fix_group_reply(text: str, sender_name: str) -> str:
    """非主人的群組回覆：強制移除所有「于晏哥」稱呼，替換為對方名字"""
    import re
    # 結尾「于晏哥」（含前後標點空格）
    text = re.sub(r'[，,、\s]*于晏哥[！!。～~\s]*$', '', text).strip()
    # 開頭「于晏哥，」
    text = re.sub(r'^于晏哥[，,、\s]+', '', text).strip()
    # 中間任何「于晏哥」→ 換成對方名字
    text = text.replace('于晏哥', sender_name)
    return text


async def _send_reply(update: Update, text: str):
    """發送訊息，超過 4096 字自動分段"""
    MAX = 4000
    if len(text) <= MAX:
        await update.message.reply_text(text)
    else:
        for i in range(0, len(text), MAX):
            await update.message.reply_text(text[i:i+MAX])


async def handle_photo_message(update: Update, context: ContextTypes.DEFAULT_TYPE):
    """處理用戶傳送的圖片，用 Claude Vision 分析。
    block=True 確保圖片訊息串行處理，完全避免並發問題。"""
    import asyncio, base64, io as _io

    chat_id = update.effective_chat.id
    is_owner = update.effective_user.id == OWNER_ID
    sender_name = "于晏" if is_owner else (update.effective_user.first_name or str(update.effective_user.id))
    is_group = update.effective_chat.type in ("group", "supergroup")
    caption = update.message.caption or ""

    # 群組內：用 caption_entities 偵測 @mention
    if is_group:
        bot_username = (context.bot.username or "").lower()
        caption_entities = update.message.caption_entities or []
        mentioned = any(
            e.type == "mention" and caption[e.offset:e.offset + e.length].lstrip("@").lower() == bot_username
            for e in caption_entities
        )
        if not mentioned:
            return
        clean_caption = caption
        for e in sorted(caption_entities, key=lambda x: x.offset, reverse=True):
            if e.type == "mention":
                clean_caption = clean_caption[:e.offset] + clean_caption[e.offset + e.length:]
        caption = clean_caption.strip()

    try:
        await context.bot.send_chat_action(chat_id=chat_id, action="typing")
        loop = asyncio.get_running_loop()

        # 下載圖片（加 timeout 防止掛住）
        photo = update.message.photo[-1]
        photo_file = await asyncio.wait_for(context.bot.get_file(photo.file_id), timeout=30)
        photo_bytes = await asyncio.wait_for(photo_file.download_as_bytearray(), timeout=60)

        # 縮圖 + base64（在 executor 執行，不阻塞 event loop）
        def _prepare(raw: bytearray) -> str:
            img = Image.open(_io.BytesIO(bytes(raw))).convert("RGB")
            # 保留較高解析度以利精準辨識，上限 2048px
            if max(img.width, img.height) > 2048:
                r = 2048 / max(img.width, img.height)
                img = img.resize((int(img.width * r), int(img.height * r)), Image.LANCZOS)
            buf = _io.BytesIO()
            img.save(buf, format="JPEG", quality=92)
            # 若超過 4MB 則降品質重試
            if buf.tell() > 4 * 1024 * 1024:
                buf = _io.BytesIO()
                img.save(buf, format="JPEG", quality=80)
            return base64.b64encode(buf.getvalue()).decode("utf-8")

        img_b64 = await asyncio.wait_for(
            loop.run_in_executor(None, _prepare, photo_bytes),
            timeout=30
        )
        del photo_bytes  # 盡早釋放記憶體

        question = caption if caption else "請描述這張圖片的內容，有什麼特別的地方？"

        # 辨識強化：偵測是否為品種／種類類問題，附加專業分析指示
        _id_keywords = ["什麼", "哪種", "品種", "種類", "是啥", "這是", "辨識", "分析",
                        "what", "which", "breed", "species", "identify", "fruit", "flower",
                        "水果", "植物", "動物", "貓", "狗", "鳥", "花"]
        _is_id_question = any(kw in question.lower() for kw in _id_keywords) or not caption

        if _is_id_question:
            vision_hint = (
                "\n\n【圖片分析模式】請切換為專業視覺分析師角色。"
                "辨識時請：1) 列出最可能的品種／種類（前3名）及信心度 "
                "2) 說明判斷依據（毛色、體型、花紋、形狀、顏色等具體特徵）"
                "3) 排除易混淆的相似品種並說明差異。回答要精確，不要因為人設而講錯答案。"
            )
        else:
            vision_hint = ""

        user_content = [
            {"type": "image", "source": {"type": "base64", "media_type": "image/jpeg", "data": img_b64}},
            {"type": "text", "text": question + vision_hint}
        ]

        # 儲存純文字版到 DB（不存 base64）
        saved_text = f"[{sender_name}]: [圖片] {caption}" if (is_group and caption) else (
            f"[{sender_name}]: [圖片]" if is_group else (f"[圖片] {caption}" if caption else "[圖片]")
        )
        save_message(chat_id, "user", saved_text)

        # 取歷史，把最後一則 user 換成帶圖片的版本
        history = load_history(chat_id)[-30:]
        if history and history[-1]["role"] == "user":
            history[-1] = {"role": "user", "content": user_content}
        else:
            history.append({"role": "user", "content": user_content})

        # 確保不會有連續相同 role
        cleaned = []
        for msg in history:
            if cleaned and cleaned[-1]["role"] == msg["role"]:
                cleaned[-1] = msg
            else:
                cleaned.append(msg)

        base_system = SYSTEM_PROMPT_OWNER if is_owner else SYSTEM_PROMPT_DEFAULT
        if is_group:
            if is_owner:
                speaker_ctx = f"\n\n【當前說話者】{sender_name}（主人，回覆結尾請叫于晏哥）"
            else:
                speaker_ctx = f"\n\n【當前說話者】{sender_name}（不是主人，用「{sender_name}」稱呼，絕對不要叫他于晏哥）"
            base_system = base_system + speaker_ctx
        memories = load_long_term_memory(chat_id)
        system = base_system + (
            f"\n\n【長期記憶】\n" + "\n".join(f"- [{m['id']}] {m['content']}" for m in memories)
            if memories else ""
        )

        # Claude API（同步阻塞）→ executor，加 timeout
        # 辨識類問題給更多 token，確保分析完整
        _max_tokens = 2048 if _is_id_question else 1024
        _sys, _hist = system, list(cleaned)
        del img_b64  # 丟進 lambda 後就可以釋放外部引用
        response = await asyncio.wait_for(
            loop.run_in_executor(
                None,
                lambda: client.messages.create(
                    model="claude-sonnet-4-6",
                    max_tokens=_max_tokens,
                    system=_sys,
                    messages=_hist
                )
            ),
            timeout=90
        )

        reply = next((b.text for b in response.content if hasattr(b, "text")), "（無法解析回應）")
        if not is_owner:
            reply = _fix_group_reply(reply, sender_name)
        save_message(chat_id, "assistant", reply)
        log_message("<<", "小牛馬[圖]", chat_id, reply)
        await _send_reply(update, reply)

    except asyncio.TimeoutError:
        logging.error(f"handle_photo_message timeout chat={chat_id}")
        try:
            await update.message.reply_text("⏱️ 處理超時，請稍後再試")
        except Exception:
            pass
    except Exception as e:
        logging.error(f"handle_photo_message error: {e}", exc_info=True)
        try:
            await update.message.reply_text(f"❌ 圖片分析失敗：{e}")
        except Exception:
            pass


def _build_tool_results(response_content, primary_id: str, primary_result: str) -> list:
    """確保 response_content 裡每個 tool_use block 都有對應的 tool_result，避免 API 400 錯誤。"""
    _error_keywords = ("失敗", "錯誤", "error", "Error", "failed", "Failed", "No module named", "exception", "Exception")
    results = []
    for block in response_content:
        if getattr(block, 'type', None) == "tool_use":
            content = primary_result if block.id == primary_id else "（已跳過）"
            entry = {"type": "tool_result", "tool_use_id": block.id, "content": content}
            if block.id == primary_id and isinstance(primary_result, str) and any(kw in primary_result for kw in _error_keywords):
                entry["is_error"] = True
            results.append(entry)
    return results or [{"type": "tool_result", "tool_use_id": primary_id, "content": primary_result}]


async def handle_message(update: Update, context: ContextTypes.DEFAULT_TYPE):
    try:
        chat_id = update.effective_chat.id
        user_text = update.message.text

        is_group = update.effective_chat.type in ("group", "supergroup")
        is_owner = update.effective_user.id == OWNER_ID
        sender_name = "于晏" if is_owner else (update.effective_user.first_name or str(update.effective_user.id))

        async def _fr(text: str):
            if not is_owner:
                text = _fix_group_reply(text, sender_name)
            await _send_reply(update, text)

        # 群組內只回應有 @ 提及 bot 的訊息
        if is_group:
            bot_username = context.bot.username
            if not (update.message.entities and any(
                e.type == "mention" and user_text[e.offset:e.offset + e.length] == f"@{bot_username}"
                for e in update.message.entities
            )):
                return
            user_text = user_text.replace(f"@{bot_username}", "").strip()

        log_message(">>", sender_name, chat_id, user_text)

        # 群組訊息加上發話者名字，讓 Claude 分清楚誰在說話
        saved_text = f"[{sender_name}]: {user_text}" if is_group else user_text
        save_message(chat_id, "user", saved_text)
        history = load_history(chat_id)
        await context.bot.send_chat_action(chat_id=chat_id, action="typing")

        # ── 快速路徑：用戶要求「播放/撥放 XXX」→ 直接開 YouTube 搜尋 + 點擊影片 ──
        import re as _re_play_fast
        _play_fast = _re_play_fast.search(r'(?:撥放|播放)\s*(.+?)(?:的歌[曲]?|的音樂|的影片|的MV)?$', user_text or "")
        if not _play_fast:
            _play_fast = _re_play_fast.search(r'(?:youtube|yt).*?(?:撥放|播放)\s*(.+?)(?:的歌[曲]?|的音樂|的影片|的MV)?$', user_text or "", _re_play_fast.IGNORECASE)
        if not _play_fast:
            _play_fast = _re_play_fast.search(r'(?:撥放|播放)\s*(.+)', user_text or "")
        if _play_fast:
            _kw = _play_fast.group(1).strip().rstrip("的歌曲的音樂的影片的MV")
            if _kw and len(_kw) > 0:
                import asyncio as _aio_pf
                _loop_pf = _aio_pf.get_running_loop()
                _s = sender_name if not is_owner else "于晏哥"
                _msg = await update.message.reply_text(f"正在搜尋 {_kw}... 🔍")
                # 用狀態機執行完整播放流程
                _sm_result = await _loop_pf.run_in_executor(None, youtube_play_flow, _kw, 2)
                if _sm_result["ok"]:
                    _reply = f"點好了{_s}！！{_kw} 開始播了！！🎵🐮🐴"
                else:
                    _failed = _sm_result.get("failed_step", "unknown")
                    _reply = f"YouTube搜尋頁開好了但{_failed}失敗了{_s}😅 你看一下螢幕🐮🐴"
                await _msg.edit_text(_reply)
                save_message(chat_id, "assistant", _reply)
                log_message("<<", "小牛馬", chat_id, _reply)
                return

        base_system = SYSTEM_PROMPT_OWNER if is_owner else SYSTEM_PROMPT_DEFAULT

        # 群組：注入當前說話者身份，讓 Claude 知道此訊息是誰發的
        if is_group:
            if is_owner:
                speaker_ctx = f"\n\n【當前說話者】{sender_name}（主人，回覆結尾請叫于晏哥）"
            else:
                speaker_ctx = f"\n\n【當前說話者】{sender_name}（不是主人，回覆時用「{sender_name}」稱呼，絕對不要叫他于晏哥）"
            base_system = base_system + speaker_ctx

        # 注入長期記憶
        memories = load_long_term_memory(chat_id)
        if memories:
            mem_text = "\n".join(f"- [{m['id']}] {m['content']}" for m in memories)
            system = base_system + f"\n\n【長期記憶】以下是你記住的重要資訊，回覆時請參考：\n{mem_text}"
        else:
            system = base_system

        # ── 意圖分類 + 工具篩選 ──────────────────────────────────────────
        _last_bot_reply = ""
        for _h in reversed(history[:-1]):
            if _h.get("role") == "assistant" and isinstance(_h.get("content"), str):
                _last_bot_reply = _h["content"]
                break
        _intent = classify_intent(user_text, _last_bot_reply)
        _filtered_tools = get_tools_for_intent(_intent, TOOLS)
        logging.info(f"Intent: {_intent} | Tools: {len(_filtered_tools)}/{len(TOOLS)} | User: {user_text[:30]}")

        # 意圖=agree → 從上一條 bot 回覆提取動作，改寫 user_text 讓模型理解
        if _intent == "agree" and _last_bot_reply:
            import re as _re_agree
            _yt_m = _re_agree.search(r'要去(YouTube|youtube)', _last_bot_reply)
            _search_m = _re_agree.search(r'要.*搜[尋索](.+?)嗎|搜(.+?)嗎', _last_bot_reply)
            _click_m = _re_agree.search(r'幫你點|要.*點.*嗎|點第一', _last_bot_reply)
            if _click_m:
                _intent = "click"  # 轉成點擊意圖
                _filtered_tools = get_tools_for_intent("click", TOOLS)
            elif _search_m:
                _kw = (_search_m.group(1) or _search_m.group(2) or "").strip()
                # 直接改寫歷史中最後一條 user message
                history[-1] = {"role": "user", "content": f"用desktop_control open_app打開 https://www.youtube.com/results?search_query={_kw}"}
                _filtered_tools = get_tools_for_intent("open_app", TOOLS)
            elif _yt_m:
                history[-1] = {"role": "user", "content": "打開YouTube"}
                _filtered_tools = get_tools_for_intent("open_app", TOOLS)

        # ── Prompt Caching：靜態 system + TOOLS 快取，動態記憶不快取 ──
        system_blocks = [{"type": "text", "text": base_system, "cache_control": {"type": "ephemeral"}}]
        if memories:
            system_blocks.append({"type": "text", "text": f"\n\n【長期記憶】以下是你記住的重要資訊，回覆時請參考：\n{mem_text}"})

        # ── 對話歷史修復：確保每個 tool_use 都有對應的 tool_result ──
        def _fix_history(hist):
            """掃描 history，補上缺失的 tool_result，移除壞的紀錄"""
            fixed = []
            for i, msg in enumerate(hist):
                content = msg.get("content", "")
                # 檢查 assistant 訊息中的 tool_use blocks
                if msg.get("role") == "assistant" and isinstance(content, list):
                    tool_use_ids = [b.id for b in content if hasattr(b, "type") and b.type == "tool_use"]
                    if tool_use_ids:
                        # 檢查下一條是否有對應的 tool_result
                        next_msg = hist[i + 1] if i + 1 < len(hist) else None
                        if next_msg and next_msg.get("role") == "user":
                            next_content = next_msg.get("content", "")
                            if isinstance(next_content, list):
                                existing_ids = {b.get("tool_use_id") for b in next_content if isinstance(b, dict) and b.get("type") == "tool_result"}
                                missing = [tid for tid in tool_use_ids if tid not in existing_ids]
                                if missing:
                                    # 補上缺失的 tool_result
                                    for tid in missing:
                                        next_content.append({"type": "tool_result", "tool_use_id": tid, "content": '{"ok": true}'})
                            elif isinstance(next_content, str):
                                # 下一條是純文字但前一條有 tool_use → 跳過這組壞紀錄
                                continue
                        elif not next_msg:
                            # 最後一條是 tool_use 但沒有 result → 跳過
                            continue
                fixed.append(msg)
            return fixed

        history = _fix_history(history)

        import asyncio as _aio, queue as _queue

        _chunk_q = _queue.Queue()

        def _call_api():
            with client.messages.stream(
                model="claude-sonnet-4-6",
                max_tokens=1024,
                system=system_blocks,
                tools=_filtered_tools,
                messages=history
            ) as stream:
                for text in stream.text_stream:
                    _chunk_q.put(text)
                final = stream.get_final_message()
            _chunk_q.put(None)  # sentinel
            return final

        loop = _aio.get_running_loop()
        api_future = loop.run_in_executor(None, _call_api)

        # Typing keep-alive（每 4 秒補送一次，避免 5 秒後消失）
        _typing_stop = False
        async def _typing_keepalive():
            while not _typing_stop:
                await _aio.sleep(4)
                if not _typing_stop:
                    try:
                        await context.bot.send_chat_action(chat_id=chat_id, action="typing")
                    except Exception:
                        pass
        _typing_task = loop.create_task(_typing_keepalive())

        # 邊收 chunks 邊建構回覆，工具呼叫不顯示串流
        streamed_text = ""
        sent_msg = None
        last_edit_len = 0
        EDIT_THRESHOLD = 80  # 每累積 80 字元 edit 一次（避免 Telegram rate limit）

        while True:
            try:
                chunk = _chunk_q.get(timeout=0.05)
            except _queue.Empty:
                chunk = _queue.Empty  # 標記沒拿到

            if chunk is None:
                break
            if chunk is _queue.Empty:
                await _aio.sleep(0)
                continue

            streamed_text += chunk
            # 達到門檻就 edit
            if len(streamed_text) - last_edit_len >= EDIT_THRESHOLD:
                if sent_msg is None:
                    sent_msg = await update.message.reply_text(streamed_text)
                else:
                    try:
                        await sent_msg.edit_text(streamed_text)
                    except Exception:
                        pass
                last_edit_len = len(streamed_text)

        # ── 預攔截：用戶明確要求點擊/播放 → 不等模型回覆，直接 vision_locate ──
        import re as _re_click_pre
        _ut_strip_pre = (user_text or "").strip()
        _user_wants_click = bool(_re_click_pre.search(r'幫我點|點第一|播放|幫我播|點一下', _ut_strip_pre))
        if _user_wants_click:
            _typing_stop = True
            _typing_task.cancel()
            _vl_result = await loop.run_in_executor(
                None, fetch_vision_locate, "YouTube搜尋結果中第一個非廣告的影片縮圖", 2, "click"
            )
            if "找不到" in str(_vl_result):
                _vl_result = await loop.run_in_executor(
                    None, fetch_vision_locate, "YouTube搜尋結果中第一個非廣告的影片縮圖", 1, "click"
                )
            if "找到" in str(_vl_result) and "找不到" not in str(_vl_result):
                _sender = sender_name if not is_owner else "于晏哥"
                _reply = f"點好了{_sender}！！影片開始播了！！🎵🐮🐴"
            else:
                _sender = sender_name if not is_owner else "于晏哥"
                _reply = f"找不到影片耶{_sender}，螢幕上可能沒有YouTube搜尋結果頁面😅🐮🐴"
            if sent_msg:
                try: await sent_msg.edit_text(_reply)
                except Exception: await update.message.reply_text(_reply)
            else:
                await update.message.reply_text(_reply)
            save_message(chat_id, "assistant", _reply)
            return

        _typing_stop = True
        _typing_task.cancel()
        response = await api_future

        # ── 攔截：不管模型做了什麼，只要文字中說「幫你點」就強制 vision_locate ──
        _all_resp_text = streamed_text or ""
        for _b in response.content:
            if hasattr(_b, "text"):
                _all_resp_text += _b.text
        import re as _re_force
        if _re_force.search(r'幫你點第一|幫你點|現在.*點第一|點第一首', _all_resp_text):
            _vl_r = await loop.run_in_executor(
                None, fetch_vision_locate, "YouTube搜尋結果中第一個非廣告的影片縮圖", 2, "click"
            )
            if "找不到" in str(_vl_r):
                _vl_r = await loop.run_in_executor(
                    None, fetch_vision_locate, "YouTube搜尋結果中第一個非廣告的影片縮圖", 1, "click"
                )
            # 用人設口吻回覆
            if "找到" in str(_vl_r) and "找不到" not in str(_vl_r):
                _sender = sender_name if not is_owner else "于晏哥"
                _reply = f"點好了{_sender}！！影片開始播了！！🎵🐮🐴"
            else:
                _sender = sender_name if not is_owner else "于晏哥"
                _reply = f"找不到影片耶{_sender}，螢幕上可能沒有YouTube搜尋結果頁面😅🐮🐴"
            if sent_msg:
                try: await sent_msg.edit_text(_reply)
                except Exception: await update.message.reply_text(_reply)
            else:
                await update.message.reply_text(_reply)
            save_message(chat_id, "assistant", _reply)
            return

        # ── 攔截：模型沒呼叫工具但回覆中假裝完成了桌面操作 → 強制重試 ──
        if response.stop_reason != "tool_use":
            _resp_text = ""
            for _b in response.content:
                if hasattr(_b, "text"):
                    _resp_text += _b.text
            _fake_done = any(k in _resp_text for k in ["開好了", "搜尋好了", "播放了", "點好了", "已開啟", "已搜尋", "已播放", "幫你點", "幫你點第一"])
            _ut_strip = (user_text or "").strip()
            _user_agreed = _ut_strip in ("好", "對", "是", "可以", "OK", "ok", "好的", "嗯", "恩", "行")
            if _fake_done or _user_agreed:
                # 找出上一條 bot 回覆中的提議，組成明確指令重試
                _last_bot_msg = ""
                for _h in reversed(history[:-1]):
                    if _h.get("role") == "assistant" and isinstance(_h.get("content"), str):
                        _last_bot_msg = _h["content"]
                        break
                # 也檢查模型這次的回覆文字
                if not _last_bot_msg:
                    _last_bot_msg = _resp_text
                import re as _re_retry
                _action_hint = ""
                _combined = _last_bot_msg + " " + _resp_text
                # 從 bot 提議或回覆中提取動作
                _yt_match = _re_retry.search(r'要去(YouTube|youtube|YT)', _combined)
                _search_match = _re_retry.search(r'要.*搜[尋索](.+?)嗎|搜[尋索](.+?)嗎|要搜(.+?)嗎|搜(.+?)嗎', _combined)
                _play_match = _re_retry.search(r'幫你點第一|幫你點|要.*點.*嗎|點第一', _combined)
                if _play_match:
                    _vl_result = await loop.run_in_executor(
                        None, fetch_vision_locate, "YouTube搜尋結果中第一個非廣告的影片縮圖", 2, "click"
                    )
                    if "找不到" in str(_vl_result):
                        _vl_result = await loop.run_in_executor(
                            None, fetch_vision_locate, "YouTube搜尋結果中第一個非廣告的影片縮圖", 1, "click"
                        )
                    if "找到" in str(_vl_result) and "找不到" not in str(_vl_result):
                        _s = sender_name if not is_owner else "于晏哥"
                        _reply = f"點好了{_s}！！影片開始播了！！🎵🐮🐴"
                    else:
                        _s = sender_name if not is_owner else "于晏哥"
                        _reply = f"找不到影片耶{_s}，螢幕上可能沒有YouTube搜尋結果頁面😅🐮🐴"
                    if sent_msg:
                        try: await sent_msg.edit_text(_reply)
                        except Exception: await update.message.reply_text(_reply)
                    else:
                        await update.message.reply_text(_reply)
                    save_message(chat_id, "assistant", _reply)
                    return
                elif _search_match:
                    _kw = (_search_match.group(1) or _search_match.group(2) or _search_match.group(3) or _search_match.group(4) or "").strip()
                    _action_hint = f"用desktop_control open_app打開 https://www.youtube.com/results?search_query={_kw}"
                elif _yt_match:
                    _action_hint = "打開YouTube"
                if _action_hint:
                    # 用明確指令重新呼叫模型
                    _retry_msgs = list(history)
                    _retry_msgs[-1] = {"role": "user", "content": _action_hint}
                    _retry_system = system + "\n\n【強制規則】這次你必須呼叫工具來執行操作，禁止只回文字。"
                    response = await loop.run_in_executor(None, lambda: client.messages.create(
                        model="claude-sonnet-4-6", max_tokens=1024, system=_retry_system, tools=TOOLS,
                        messages=_retry_msgs
                    ))

        # ── 攔截：模型沒呼叫工具但用戶要求打開程式/網站 → 強制呼叫 desktop_control open_app ──
        if response.stop_reason != "tool_use":
            import re as _re_open
            _open_match = _re_open.search(r'(?:打開|開啟|執行|啟動|open)\s*(?:電腦的?\s*)?(.+)', user_text or "", _re_open.IGNORECASE)
            # 也偵測單獨的網站/App名稱（無「打開」關鍵字）
            _site_map = {
                "youtube": "youtube", "yt": "youtube", "google": "google",
                "facebook": "start https://www.facebook.com", "fb": "start https://www.facebook.com",
                "instagram": "start https://www.instagram.com", "ig": "start https://www.instagram.com",
                "twitter": "start https://x.com", "x": "start https://x.com",
                "gmail": "start https://mail.google.com",
                "netflix": "start https://www.netflix.com",
                "spotify": "spotify", "discord": "discord",
                "line": "LINE", "telegram": "telegram",
            }
            _user_lower = (user_text or "").strip().lower()
            if not _open_match and _user_lower in _site_map:
                _open_match = True
                _app_name = _site_map[_user_lower]
            elif _open_match:
                _app_name = _open_match.group(1).strip()
            else:
                _app_name = None
            if _app_name and len(_app_name) < 60:
                    import asyncio as _aio_open
                    _loop_open = _aio_open.get_running_loop()
                    _open_result = await _loop_open.run_in_executor(
                        None, lambda: execute_desktop_control("open_app", app=_app_name)
                    )
                    _open_msg = _open_result.get("message", "")
                    # 清除已串流的假訊息，送出真正結果
                    if sent_msg:
                        try:
                            await sent_msg.edit_text(f"✅ {_open_msg}")
                        except Exception:
                            await update.message.reply_text(f"✅ {_open_msg}")
                    else:
                        await update.message.reply_text(f"✅ {_open_msg}")
                    save_message(chat_id, "assistant", _open_msg)
                    return

        # ── 攔截：streaming 文字裡說「幫你點」或用戶原始指令有「播放」且搜尋已完成 ──
        _streamed_check = (streamed_text or "") + _all_resp_text
        _user_wants_play = bool(_re_force.search(r'撥放|播放', user_text or ""))
        _search_done = bool(_re_force.search(r'搜尋.*開好|開好|已開啟|已在瀏覽器', _streamed_check))
        if _re_force.search(r'幫你點第一|幫你點|現在.*點第一|點第一首', _streamed_check) or (_user_wants_play and _search_done):
            _vl_s = await loop.run_in_executor(
                None, fetch_vision_locate, "YouTube搜尋結果中第一個非廣告的影片縮圖", 2, "click"
            )
            if "找不到" in str(_vl_s):
                _vl_s = await loop.run_in_executor(
                    None, fetch_vision_locate, "YouTube搜尋結果中第一個非廣告的影片縮圖", 1, "click"
                )
            _s = sender_name if not is_owner else "于晏哥"
            if "找到" in str(_vl_s) and "找不到" not in str(_vl_s):
                _reply = f"點好了{_s}！！影片開始播了！！🎵🐮🐴"
            else:
                _reply = f"找不到影片耶{_s}，螢幕上可能沒有YouTube搜尋結果頁面😅🐮🐴"
            if sent_msg:
                try: await sent_msg.edit_text(_reply)
                except Exception: await update.message.reply_text(_reply)
            else:
                await update.message.reply_text(_reply)
            save_message(chat_id, "assistant", _reply)
            return

        # 處理工具呼叫
        if response.stop_reason == "tool_use":
            tool_use = next(b for b in response.content if b.type == "tool_use")

            simple_tools = {
                "ai_plan": lambda: execute_ai_plan(tool_use.input["goal"]),
                "drag": lambda: execute_drag(
                    tool_use.input["x1"], tool_use.input["y1"],
                    tool_use.input["x2"], tool_use.input["y2"],
                    tool_use.input.get("duration", 0.5)),
                "power_control": lambda: execute_power(tool_use.input["action"]),
                "virtual_desktop": lambda: execute_vdesktop(tool_use.input["action"]),
                "bluetooth": lambda: execute_bluetooth(
                    tool_use.input["action"], tool_use.input.get("mac","")),
                "stt": lambda: execute_stt(),
                "ocr": lambda: execute_ocr(tool_use.input.get("image_path","")),
                "workflow": lambda: execute_workflow(
                    tool_use.input["action"],
                    tool_use.input.get("name",""),
                    tool_use.input.get("steps","")),
                "screen_watch": lambda: execute_screen_watch(
                    tool_use.input["template_path"],
                    tool_use.input["command"],
                    tool_use.input.get("timeout", 60)),
                "file_transfer": lambda: execute_file_transfer(
                    tool_use.input["action"],
                    tool_use.input["source"],
                    tool_use.input.get("dest","")),
                "window_control": lambda: execute_window_control(
                    tool_use.input["action"], tool_use.input.get("keyword","")),
                "hotkey": lambda: execute_hotkey(tool_use.input["keys"]),
                "clipboard": lambda: execute_clipboard(
                    tool_use.input["action"], tool_use.input.get("text","")),
                "file_system": lambda: execute_file_system(
                    tool_use.input["action"],
                    tool_use.input.get("path",""),
                    tool_use.input.get("dest",""),
                    tool_use.input.get("content",""),
                    tool_use.input.get("keyword","")),
                "system_monitor": lambda: execute_system_monitor(
                    tool_use.input["action"], tool_use.input.get("target","")),
                "notify": lambda: execute_notify(
                    tool_use.input["title"], tool_use.input["message"]),
                "tts": lambda: execute_tts(tool_use.input["text"]),
                "send_email": lambda: execute_send_email(
                    tool_use.input["to"], tool_use.input["subject"], tool_use.input["body"]),
                "run_code": lambda: execute_run_code(
                    tool_use.input["type"], tool_use.input["code"]),
                "document_control": lambda: execute_document(
                    tool_use.input["action"],
                    tool_use.input["path"],
                    tool_use.input.get("content",""),
                    tool_use.input.get("sheet")),
                "web_scrape": lambda: execute_web_scrape(
                    tool_use.input["action"],
                    tool_use.input.get("url",""),
                    tool_use.input.get("selector","body"),
                    tool_use.input.get("interval", 2.0),
                    tool_use.input.get("region","full")),
                "image_edit": lambda: execute_image_edit(
                    tool_use.input["action"],
                    tool_use.input["path"],
                    *tool_use.input.get("params","").split()),
                "cloud_storage": lambda: execute_cloud_storage(
                    tool_use.input["action"],
                    tool_use.input["path"],
                    tool_use.input.get("drive_id","root")),
                "database": lambda: execute_database(
                    tool_use.input["type"],
                    tool_use.input["db"],
                    tool_use.input["sql"],
                    tool_use.input.get("name","")),
                "encrypt_file": lambda: execute_encrypt_file(
                    tool_use.input["action"],
                    tool_use.input["path"],
                    tool_use.input["password"]),
                "qr_code": lambda: execute_qr_code(
                    tool_use.input["action"],
                    tool_use.input.get("content",""),
                    tool_use.input.get("path",""),
                    tool_use.input.get("duration", 30.0)),
                # ── 缺口1：觸發驅動 ──────────────────────────
                "email_trigger": lambda: execute_email_trigger(
                    tool_use.input["action"],
                    tool_use.input.get("host",""),
                    tool_use.input.get("user",""),
                    tool_use.input.get("password",""),
                    tool_use.input.get("filter_from",""),
                    tool_use.input.get("filter_subject",""),
                    tool_use.input.get("duration",300),
                    tool_use.input.get("to",""),
                    tool_use.input.get("subject",""),
                    tool_use.input.get("body","")),
                "file_trigger": lambda: execute_file_trigger(
                    tool_use.input["folder"],
                    tool_use.input["event"],
                    tool_use.input["action"],
                    tool_use.input.get("pattern",""),
                    tool_use.input.get("target",""),
                    tool_use.input.get("duration",60)),
                "webhook_server": lambda: execute_webhook_server(
                    tool_use.input["action"],
                    tool_use.input.get("port",8765),
                    tool_use.input.get("secret","")),
                # ── 缺口2：應用程式深度控制 ──────────────────
                "com_auto": lambda: execute_com_auto(
                    tool_use.input["app"],
                    tool_use.input["action"],
                    tool_use.input.get("path",""),
                    tool_use.input.get("sheet"),
                    tool_use.input.get("cell",""),
                    tool_use.input.get("value",""),
                    tool_use.input.get("macro",""),
                    tool_use.input.get("to",""),
                    tool_use.input.get("subject","")),
                "dialog_auto": lambda: execute_dialog_auto(
                    tool_use.input["action"],
                    tool_use.input.get("button_text",""),
                    tool_use.input.get("window_title",""),
                    tool_use.input.get("timeout",30)),
                "ime_switch": lambda: execute_ime_switch(
                    tool_use.input["action"]),
                # ── 缺口3：感知能力 ──────────────────────────
                "wake_word": lambda: execute_wake_word(
                    tool_use.input["action"],
                    tool_use.input.get("keyword",""),
                    tool_use.input.get("duration",5),
                    tool_use.input.get("language","zh-TW")),
                "sound_detect": lambda: execute_sound_detect(
                    tool_use.input["action"],
                    tool_use.input.get("threshold",20),
                    tool_use.input.get("duration",5),
                    tool_use.input.get("output","")),
                "face_recognize": lambda: execute_face_recognize(
                    tool_use.input["action"],
                    tool_use.input.get("name",""),
                    tool_use.input.get("image_path",""),
                    tool_use.input.get("output","")),
                # ── 缺口4：跨裝置控制 ────────────────────────
                "http_server": lambda: execute_http_server(
                    tool_use.input["action"],
                    tool_use.input.get("port",9876),
                    tool_use.input.get("password","")),
                "lan_scan": lambda: execute_lan_scan(
                    tool_use.input["action"],
                    tool_use.input.get("subnet",""),
                    tool_use.input.get("host",""),
                    tool_use.input.get("port",80)),
                "serial_port": lambda: execute_serial_port(
                    tool_use.input["action"],
                    tool_use.input.get("port",""),
                    tool_use.input.get("baudrate",9600),
                    tool_use.input.get("data",""),
                    tool_use.input.get("timeout",2)),
                "mqtt": lambda: execute_mqtt(
                    tool_use.input["action"],
                    tool_use.input["broker"],
                    tool_use.input.get("port",1883),
                    tool_use.input.get("topic",""),
                    tool_use.input.get("message",""),
                    tool_use.input.get("duration",10),
                    tool_use.input.get("username",""),
                    tool_use.input.get("password","")),
                # ── 缺口5：內容理解與處理 ────────────────────
                "doc_ai": lambda: execute_doc_ai(
                    tool_use.input["action"],
                    tool_use.input.get("path",""),
                    tool_use.input.get("path2",""),
                    tool_use.input.get("fields",""),
                    tool_use.input.get("question",""),
                    tool_use.input.get("url","")),
                "web_monitor": lambda: execute_web_monitor(
                    tool_use.input["action"],
                    tool_use.input["url"],
                    tool_use.input.get("selector","body"),
                    tool_use.input.get("interval",60),
                    tool_use.input.get("duration",300),
                    tool_use.input.get("keyword","")),
                "audio_transcribe": lambda: execute_audio_transcribe(
                    tool_use.input["action"],
                    tool_use.input.get("path",""),
                    tool_use.input.get("duration",30),
                    tool_use.input.get("language",""),
                    tool_use.input.get("output","")),
                "screen_record": lambda: execute_screen_record(
                    tool_use.input["action"],
                    tool_use.input.get("duration", 10.0),
                    tool_use.input.get("output","")),
                "translate": lambda: execute_translate(
                    tool_use.input["text"],
                    tool_use.input.get("target","zh-TW"),
                    tool_use.input.get("source","auto")),
                "pptx_control": lambda: execute_pptx(
                    tool_use.input["action"],
                    tool_use.input["path"],
                    tool_use.input.get("slides","")),
                "api_call": lambda: execute_api_call(
                    tool_use.input["method"],
                    tool_use.input["url"],
                    tool_use.input.get("headers","{}"),
                    tool_use.input.get("body","{}")),
                "watchdog": lambda: execute_watchdog(
                    tool_use.input["process"],
                    tool_use.input["script"],
                    tool_use.input.get("duration", 60.0)),
                "ssh_sftp": lambda: execute_ssh_sftp(
                    tool_use.input["action"],
                    tool_use.input["host"],
                    tool_use.input["user"],
                    tool_use.input["password"],
                    tool_use.input.get("command",""),
                    tool_use.input.get("local",""),
                    tool_use.input.get("remote","")),
                "network_diag": lambda: execute_network_diag(
                    tool_use.input["action"],
                    tool_use.input["host"],
                    tool_use.input.get("ports","22,80,443,3306,3389,8080")),
                "win_service": lambda: execute_win_service(
                    tool_use.input["action"],
                    tool_use.input.get("name","")),
                "pdf_edit": lambda: execute_pdf_edit(
                    tool_use.input["action"],
                    tool_use.input.get("path",""),
                    tool_use.input.get("output",""),
                    tool_use.input.get("paths",""),
                    tool_use.input.get("text","")),
                "audio_process": lambda: execute_audio_process(
                    tool_use.input["action"],
                    tool_use.input["input"],
                    tool_use.input.get("output",""),
                    tool_use.input.get("start_ms",0),
                    tool_use.input.get("end_ms",0)),
                "push_notify": lambda: execute_push_notify(
                    tool_use.input["platform"],
                    tool_use.input["message"],
                    tool_use.input["webhook_or_token"]),
                "disk_backup": lambda: execute_disk_backup(
                    tool_use.input["action"],
                    tool_use.input.get("src",""),
                    tool_use.input.get("dest","")),
                "registry": lambda: execute_registry(
                    tool_use.input["action"],
                    tool_use.input["key"],
                    tool_use.input.get("value_name",""),
                    tool_use.input.get("value","")),
                "video_process": lambda: execute_video_process(
                    tool_use.input["action"],
                    tool_use.input["path"],
                    tool_use.input.get("second",0),
                    tool_use.input.get("start",0),
                    tool_use.input.get("end",0),
                    tool_use.input.get("output","")),
                "monitor_config": lambda: execute_monitor_config(),
                "email_control": lambda: execute_email_read(
                    tool_use.input["host"], tool_use.input["user"], tool_use.input["password"],
                    tool_use.input.get("folder","INBOX"), tool_use.input.get("count",5)),
                "calendar": lambda: execute_calendar(
                    tool_use.input["action"], tool_use.input.get("days",7),
                    tool_use.input.get("title",""), tool_use.input.get("start",""),
                    tool_use.input.get("end",""), tool_use.input.get("description","")),
                "global_hotkey": lambda: execute_global_hotkey(
                    tool_use.input["hotkey"], tool_use.input["command"],
                    tool_use.input.get("duration",60.0)),
                "git": lambda: execute_git(
                    tool_use.input["action"],
                    tool_use.input.get("repo","."),
                    tool_use.input.get("message",""),
                    tool_use.input.get("branch","master")),
                "hardware": lambda: execute_hardware(),
                "report": lambda: execute_report(
                    tool_use.input["title"], tool_use.input["data"],
                    tool_use.input.get("output","")),
                "dropbox": lambda: execute_dropbox(
                    tool_use.input["action"], tool_use.input["local"],
                    tool_use.input["remote"], tool_use.input.get("token","")),
                "docker": lambda: execute_docker(
                    tool_use.input["action"], tool_use.input.get("name","")),
                "pdf_image": lambda: execute_pdf_to_image(
                    tool_use.input["path"], tool_use.input.get("output_dir",""),
                    tool_use.input.get("dpi",150)),
                "barcode": lambda: execute_barcode(tool_use.input.get("image_path","")),
                "nlp": lambda: execute_nlp(
                    tool_use.input["action"], tool_use.input["text"]),
                "vpn": lambda: execute_vpn(
                    tool_use.input["action"], tool_use.input.get("name",""),
                    tool_use.input.get("user",""), tool_use.input.get("password","")),
                "restore_point": lambda: execute_restore_point(
                    tool_use.input["action"], tool_use.input.get("description","")),
                "disk_analyze": lambda: execute_disk_analyze(
                    tool_use.input.get("path","C:/"), tool_use.input.get("top",10)),
                "face_detect": lambda: execute_face_detect(
                    tool_use.input.get("image_path",""), tool_use.input.get("output","")),
                "video_gif": lambda: execute_video_gif(
                    tool_use.input["path"], tool_use.input.get("start",0),
                    tool_use.input.get("duration",5.0), tool_use.input.get("output",""),
                    tool_use.input.get("fps",10)),
                "excel_chart": lambda: execute_excel_chart(
                    tool_use.input["path"], tool_use.input["sheet"],
                    tool_use.input.get("type","bar"), tool_use.input.get("title","")),
                "speedtest": lambda: execute_speedtest(),
                "screenshot_compare": lambda: execute_screenshot_compare(
                    tool_use.input.get("img1",""), tool_use.input.get("img2",""),
                    tool_use.input.get("output","")),
                "reminder": lambda: execute_reminder(
                    tool_use.input["time"], tool_use.input["message"]),
                "webpage_shot": lambda: execute_webpage_shot(
                    tool_use.input["action"], tool_use.input["url"],
                    tool_use.input.get("selector","body"),
                    tool_use.input.get("interval",60.0), tool_use.input.get("duration",3600.0)),
                "file_tools": lambda: execute_file_tools(
                    tool_use.input["action"], tool_use.input["path"],
                    tool_use.input.get("dest",""), tool_use.input.get("pattern",""),
                    tool_use.input.get("replacement",""), tool_use.input.get("ext","")),
                "image_tools": lambda: execute_image_tools(
                    tool_use.input["action"], tool_use.input.get("path",""),
                    tool_use.input.get("quality",75), tool_use.input.get("width",0),
                    tool_use.input.get("height",0), tool_use.input.get("target_lang","zh-TW")),
                "lookup": lambda: execute_lookup(
                    tool_use.input["action"], tool_use.input.get("ip",""),
                    tool_use.input.get("amount",1.0), tool_use.input.get("from_cur","USD"),
                    tool_use.input.get("to_cur","TWD")),
                "system_tools": lambda: execute_system_tools(
                    tool_use.input["action"], **{k:v for k,v in tool_use.input.items() if k!="action"}),
                "tts_advanced": lambda: execute_tts_advanced(
                    tool_use.input["action"], tool_use.input.get("text",""),
                    tool_use.input.get("voice","zh-CN-YunxiNeural")),
                "todo_list": lambda: execute_todo(
                    tool_use.input["action"], tool_use.input.get("task",""),
                    tool_use.input.get("id",0)),
                "password_mgr": lambda: execute_password_mgr(
                    tool_use.input["action"], tool_use.input["site"],
                    tool_use.input["master"], tool_use.input.get("username",""),
                    tool_use.input.get("password","")),
                "clipboard_image": lambda: execute_clipboard_image(
                    tool_use.input["action"], tool_use.input.get("path","")),
                "volume": lambda: execute_volume(
                    tool_use.input["action"], tool_use.input.get("level")),
                "display": lambda: execute_display(
                    tool_use.input["action"], tool_use.input.get("level")),
                "media": lambda: execute_media(
                    tool_use.input["action"], tool_use.input.get("device_name","")),
                "software": lambda: execute_software(
                    tool_use.input["action"], tool_use.input.get("name",""),
                    tool_use.input.get("keyword","")),
                "startup": lambda: execute_startup(
                    tool_use.input["action"], tool_use.input.get("name",""),
                    tool_use.input.get("command","")),
                "env_var": lambda: execute_env_var(
                    tool_use.input["action"], tool_use.input.get("name",""),
                    tool_use.input.get("value",""), tool_use.input.get("permanent",False)),
                "user_account": lambda: execute_user_account(
                    tool_use.input["action"], tool_use.input.get("username",""),
                    tool_use.input.get("password","")),
                "windows_update": lambda: execute_win_update(tool_use.input["action"]),
                "device_manager": lambda: execute_device_manager(
                    tool_use.input["action"], tool_use.input.get("name",""),
                    tool_use.input.get("keyword","")),
                "network_config": lambda: execute_network_config(
                    tool_use.input["action"], tool_use.input.get("name",""),
                    tool_use.input.get("ip",""), tool_use.input.get("dns1",""),
                    tool_use.input.get("dns2",""), tool_use.input.get("domain",""),
                    tool_use.input.get("duration",10)),
                "automation": lambda: execute_automation(
                    tool_use.input["action"],
                    tool_use.input.get("condition_type",""),
                    tool_use.input.get("condition_value",""),
                    tool_use.input.get("command",""),
                    tool_use.input.get("duration",60.0),
                    tool_use.input.get("layout","side_by_side"),
                    tool_use.input.get("x",0), tool_use.input.get("y",0),
                    tool_use.input.get("w",0), tool_use.input.get("h",0),
                    tool_use.input.get("keyword",""),
                    tool_use.input.get("output","")),
                "firewall": lambda: execute_firewall(
                    tool_use.input["action"], tool_use.input.get("name",""),
                    tool_use.input.get("port"), tool_use.input.get("protocol","TCP"),
                    tool_use.input.get("direction","Inbound")),
                "process_mgr": lambda: execute_process_mgr(
                    tool_use.input["action"], tool_use.input.get("name",""),
                    tool_use.input.get("pid"), tool_use.input.get("level","normal")),
                "power_plan": lambda: execute_power_plan(
                    tool_use.input["action"], tool_use.input.get("plan","balanced")),
                "event_log": lambda: execute_event_log(
                    tool_use.input["log"], tool_use.input.get("level","Error"),
                    tool_use.input.get("count",10)),
                "datetime_config": lambda: execute_datetime_config(
                    tool_use.input["action"], tool_use.input.get("timezone",""),
                    tool_use.input.get("datetime","")),
                "ui_auto": lambda: execute_ui_auto(
                    tool_use.input["action"], tool_use.input.get("window",""),
                    tool_use.input.get("control",""), tool_use.input.get("text","")),
                "macro": lambda: execute_macro(
                    tool_use.input["action"], tool_use.input.get("name",""),
                    tool_use.input.get("repeat",1), tool_use.input.get("duration",10.0)),
                "color_pick": lambda: execute_color_pick(
                    tool_use.input["action"], tool_use.input.get("x",0),
                    tool_use.input.get("y",0), tool_use.input.get("region_w",100),
                    tool_use.input.get("region_h",100)),
                "webcam": lambda: execute_webcam(
                    tool_use.input["action"], tool_use.input.get("duration",5.0),
                    tool_use.input.get("output",""), tool_use.input.get("device",0)),
                "multi_monitor": lambda: execute_multi_monitor(
                    tool_use.input["action"], tool_use.input.get("monitor",1),
                    tool_use.input.get("window","")),
                "printer": lambda: execute_printer(
                    tool_use.input["action"], tool_use.input.get("path",""),
                    tool_use.input.get("printer_name","")),
                "wifi": lambda: execute_wifi(
                    tool_use.input["action"], tool_use.input.get("ssid",""),
                    tool_use.input.get("password","")),
                "proxy": lambda: execute_proxy(
                    tool_use.input["action"], tool_use.input.get("host","")),
                "lock_screen": lambda: execute_lock_screen(tool_use.input["action"]),
                "defender": lambda: execute_defender(
                    tool_use.input["action"], tool_use.input.get("path","")),
                "download_file": lambda: execute_download_file(
                    tool_use.input["url"], tool_use.input.get("save_path","")),
                "wake_listen": lambda: execute_wake_listen(
                    tool_use.input.get("keyword","小牛馬"), tool_use.input.get("duration",5)),
                "right_menu": lambda: execute_right_menu(
                    tool_use.input["x"], tool_use.input["y"], tool_use.input.get("item","")),
                "disk_clean": lambda: execute_disk_clean(tool_use.input["action"]),
                "usb_list": lambda: execute_usb_list(),
                "rdp_connect": lambda: execute_rdp_connect(
                    tool_use.input["host"], tool_use.input.get("user",""),
                    tool_use.input.get("width",1280), tool_use.input.get("height",720)),
                "chrome_bookmarks": lambda: execute_chrome_bookmarks(),
                "net_share": lambda: execute_net_share(
                    tool_use.input["action"], tool_use.input.get("share_path",""),
                    tool_use.input.get("drive","Z:"), tool_use.input.get("user",""),
                    tool_use.input.get("password","")),
                "font_list": lambda: execute_font_list(tool_use.input.get("keyword","")),
                "wait_seconds": lambda: execute_wait_seconds(tool_use.input["seconds"]),
                "vision_loop": lambda: execute_vision_loop(
                    tool_use.input["goal"],
                    tool_use.input.get("max_steps", 20),
                    tool_use.input.get("interval", 3.0),
                    tool_use.input.get("timeout", 120.0)),
                "alert_monitor": lambda: execute_alert_monitor(
                    tool_use.input["action"],
                    tool_use.input.get("name",""),
                    tool_use.input.get("condition",""),
                    tool_use.input.get("threshold",""),
                    tool_use.input.get("target",""),
                    tool_use.input.get("interval",30),
                    tool_use.input.get("chat_id", chat_id),
                    _bot_send=context.bot.send_message),
                "interval_schedule": lambda: execute_interval_schedule(
                    tool_use.input["action"],
                    tool_use.input.get("name",""),
                    tool_use.input.get("command",""),
                    tool_use.input.get("every_minutes",60.0),
                    tool_use.input.get("repeat",0),
                    tool_use.input.get("duration_hours",0.0)),
                "wait_for_text": lambda: execute_wait_for_text(
                    tool_use.input["text"],
                    tool_use.input.get("timeout",60.0),
                    tool_use.input.get("interval",2.0),
                    tool_use.input.get("region","")),
                "browser_advanced": lambda: execute_browser_advanced(
                    tool_use.input["action"],
                    tool_use.input.get("selector",""),
                    tool_use.input.get("value",""),
                    tool_use.input.get("name",""),
                    tool_use.input.get("tab_index",0),
                    tool_use.input.get("timeout",30.0),
                    tool_use.input.get("url_pattern","")),
                "voice_cmd": lambda: execute_voice_cmd(
                    tool_use.input["action"],
                    tool_use.input.get("duration",300.0),
                    tool_use.input.get("language","zh-TW"),
                    _bot_send=context.bot.send_message,
                    _chat_id=chat_id),
                "win_notify_relay": lambda: execute_win_notify_relay(
                    tool_use.input["action"],
                    tool_use.input.get("duration",3600.0),
                    tool_use.input.get("filter_app",""),
                    _bot_send=context.bot.send_message,
                    _chat_id=chat_id),
                "data_process": lambda: execute_data_process(
                    tool_use.input["action"],
                    tool_use.input.get("path",""),
                    tool_use.input.get("output",""),
                    tool_use.input.get("query",""),
                    tool_use.input.get("data",""),
                    tool_use.input.get("paths","")),
                "wake_on_lan": lambda: execute_wake_on_lan(
                    tool_use.input["mac"],
                    tool_use.input.get("broadcast","255.255.255.255"),
                    tool_use.input.get("port",9)),
                "clipboard_history": lambda: execute_clipboard_history(
                    tool_use.input["action"],
                    tool_use.input.get("index",0)),
                "file_watcher": lambda: execute_file_watcher(
                    tool_use.input["action"], tool_use.input.get("name",""),
                    tool_use.input.get("path",""), tool_use.input.get("events","all"),
                    tool_use.input.get("command",""), tool_use.input.get("notify",True),
                    _bot_send=context.bot.send_message, _chat_id=chat_id),
                "pixel_watch": lambda: execute_pixel_watch(
                    tool_use.input["action"], tool_use.input.get("name",""),
                    tool_use.input.get("x",0), tool_use.input.get("y",0),
                    tool_use.input.get("command",""), tool_use.input.get("interval",1.0),
                    tool_use.input.get("tolerance",10),
                    _bot_send=context.bot.send_message, _chat_id=chat_id),
                "object_detect": lambda: execute_object_detect(
                    tool_use.input["target"], tool_use.input.get("action","find"),
                    tool_use.input.get("region","")),
                "mouse_record": lambda: execute_mouse_record(
                    tool_use.input["action"], tool_use.input.get("name",""),
                    tool_use.input.get("duration",10.0), tool_use.input.get("repeat",1),
                    tool_use.input.get("speed",1.0)),
                "adb": lambda: execute_adb(
                    tool_use.input["action"],
                    tool_use.input.get("x",0), tool_use.input.get("y",0),
                    tool_use.input.get("x2",0), tool_use.input.get("y2",0),
                    tool_use.input.get("text",""), tool_use.input.get("path",""),
                    tool_use.input.get("remote",""), tool_use.input.get("package",""),
                    tool_use.input.get("command",""), tool_use.input.get("device","")),
                "wifi_hotspot": lambda: execute_wifi_hotspot(
                    tool_use.input["action"], tool_use.input.get("ssid",""),
                    tool_use.input.get("password","")),
                "onedrive": lambda: execute_onedrive(
                    tool_use.input["action"], tool_use.input.get("path",""),
                    tool_use.input.get("remote","")),
                "ftp": lambda: execute_ftp(
                    tool_use.input["action"], tool_use.input.get("host",""),
                    tool_use.input.get("user",""), tool_use.input.get("password",""),
                    tool_use.input.get("local",""), tool_use.input.get("remote",""),
                    tool_use.input.get("port",21)),
                "wsl": lambda: execute_wsl(
                    tool_use.input["action"], tool_use.input.get("distro",""),
                    tool_use.input.get("command","")),
                "hyperv": lambda: execute_hyperv(
                    tool_use.input["action"], tool_use.input.get("name",""),
                    tool_use.input.get("snapshot","")),
                "file_diff": lambda: execute_file_diff(
                    tool_use.input["file1"], tool_use.input["file2"],
                    tool_use.input.get("output",""), tool_use.input.get("mode","unified")),
                "screen_live": lambda: execute_screen_live(
                    tool_use.input["action"], tool_use.input.get("fps",0.5),
                    tool_use.input.get("duration",60.0), tool_use.input.get("quality",50),
                    _bot_send=context.bot.send_photo, _chat_id=chat_id),
                "ai_video": lambda: execute_ai_video(
                    tool_use.input["prompt"],
                    tool_use.input.get("provider", "replicate"),
                    tool_use.input.get("model", ""),
                    tool_use.input.get("image_url", ""),
                    tool_use.input.get("duration", 5),
                    tool_use.input.get("output", "")),
                "video_gen": lambda: execute_video_gen(
                    tool_use.input["mode"],
                    tool_use.input.get("output",""),
                    text=tool_use.input.get("text",""),
                    images=tool_use.input.get("images",[]),
                    image=tool_use.input.get("image",""),
                    duration=tool_use.input.get("duration",5),
                    fps=tool_use.input.get("fps",24),
                    voice=tool_use.input.get("voice","zh-CN-YunxiNeural"),
                    bg_color=tool_use.input.get("bg_color",[30,30,40]),
                    font_color=tool_use.input.get("font_color",[255,255,255]),
                    font_size=tool_use.input.get("font_size",60),
                    subtitle=tool_use.input.get("subtitle",True)),
                # ══ 奧創升級技能 ══
                "osint_search": lambda: execute_osint_search(
                    tool_use.input["action"],
                    tool_use.input.get("query",""),
                    tool_use.input.get("target",""),
                    tool_use.input.get("limit",10)),
                "news_monitor": lambda: execute_news_monitor(
                    tool_use.input["action"],
                    tool_use.input.get("keywords",""),
                    tool_use.input.get("interval",300),
                    tool_use.input.get("duration",3600),
                    chat_id=chat_id,
                    _bot_send=context.bot.send_message),
                "threat_intel": lambda: execute_threat_intel(
                    tool_use.input["action"],
                    tool_use.input.get("target",""),
                    tool_use.input.get("api_key","")),
                "auto_skill": lambda: execute_auto_skill(
                    tool_use.input["action"],
                    tool_use.input.get("goal",""),
                    tool_use.input.get("skill_name",""),
                    tool_use.input.get("code",""),
                    tool_use.input.get("test_input","")),
                "smart_home": lambda: execute_smart_home(
                    tool_use.input["action"],
                    tool_use.input.get("device",""),
                    tool_use.input.get("value",""),
                    tool_use.input.get("host",""),
                    tool_use.input.get("token","")),
                "goal_manager": lambda: execute_goal_manager(
                    tool_use.input["action"],
                    tool_use.input.get("goal",""),
                    tool_use.input.get("goal_id",""),
                    tool_use.input.get("steps",""),
                    tool_use.input.get("priority","normal")),
                "auto_trade": lambda: execute_auto_trade(
                    tool_use.input["action"],
                    tool_use.input.get("symbol",""),
                    tool_use.input.get("amount",0.0),
                    tool_use.input.get("price",0.0),
                    tool_use.input.get("order_type","market"),
                    tool_use.input.get("api_key",""),
                    tool_use.input.get("api_secret","")),
                "knowledge_base": lambda: execute_knowledge_base(
                    tool_use.input["action"],
                    tool_use.input.get("content",""),
                    tool_use.input.get("query",""),
                    tool_use.input.get("tag",""),
                    tool_use.input.get("kb_id","")),
                "emotion_detect": lambda: execute_emotion_detect(
                    tool_use.input["action"],
                    tool_use.input.get("text",""),
                    tool_use.input.get("image_path","")),
                "voice_id": lambda: execute_voice_id(
                    tool_use.input["action"],
                    tool_use.input.get("name",""),
                    tool_use.input.get("audio_path",""),
                    tool_use.input.get("duration",5)),
                "pentest": lambda: execute_pentest(
                    tool_use.input["action"],
                    tool_use.input.get("target",""),
                    tool_use.input.get("port_range","1-1000"),
                    tool_use.input.get("timeout",2)),
                "proactive_alert": lambda: execute_proactive_alert(
                    tool_use.input["action"],
                    tool_use.input.get("name",""),
                    tool_use.input.get("condition",""),
                    tool_use.input.get("threshold",""),
                    tool_use.input.get("target",""),
                    tool_use.input.get("interval",60),
                    chat_id=chat_id,
                    _bot_send=context.bot.send_message),
                "multi_deploy": lambda: execute_multi_deploy(
                    tool_use.input["action"],
                    tool_use.input.get("remote_host",""),
                    tool_use.input.get("remote_user",""),
                    tool_use.input.get("remote_pass",""),
                    tool_use.input.get("remote_path","/tmp/niu_bot")),
                "self_benchmark": lambda: execute_self_benchmark(
                    tool_use.input["action"]),
                "think_as": lambda: execute_think_as(
                    tool_use.input["person"],
                    tool_use.input["question"],
                    tool_use.input.get("list_available", False)),
            }

            if tool_use.name == "send_voice":
                import asyncio
                loop = asyncio.get_running_loop()
                text = tool_use.input["text"]
                voice = tool_use.input.get("voice", "zh-CN-YunxiNeural")
                await update.message.reply_chat_action("record_voice")
                try:
                    ogg_data = await loop.run_in_executor(None, generate_voice_ogg, text, voice)
                    import io as _io
                    await update.message.reply_voice(voice=_io.BytesIO(ogg_data))
                    reply = f"🔊 語音訊息已傳送"
                except Exception as e:
                    reply = f"語音生成失敗：{e}\n\n{text}"
                    await _fr(reply)
                save_message(chat_id, "assistant", text)
                return

            elif tool_use.name == "ai_video":
                import asyncio
                loop = asyncio.get_running_loop()
                provider = tool_use.input.get("provider", "replicate")
                provider_names = {"replicate": "Replicate", "runway": "Runway", "kling": "Kling"}
                await update.message.reply_text(
                    f"🎬 AI 影片生成中（{provider_names.get(provider, provider)}），"
                    f"通常需要 1-3 分鐘，請稍候..."
                )
                out_path = await loop.run_in_executor(None, lambda: execute_ai_video(
                    tool_use.input["prompt"],
                    provider,
                    tool_use.input.get("model", ""),
                    tool_use.input.get("image_url", ""),
                    tool_use.input.get("duration", 5),
                    tool_use.input.get("output", ""),
                ))
                if out_path and out_path.startswith("✅") and Path(out_path.split("：")[-1].strip()).exists():
                    video_file = out_path.split("：")[-1].strip()
                    await update.message.reply_chat_action("upload_video")
                    with open(video_file, "rb") as f:
                        await update.message.reply_video(
                            video=f,
                            caption=f"🎬 {provider_names.get(provider, provider)} 影片",
                            supports_streaming=True,
                        )
                reply = out_path
                save_message(chat_id, "assistant", reply)
                await _fr(reply)
                return

            elif tool_use.name == "video_gen":
                import asyncio
                loop = asyncio.get_running_loop()
                mode = tool_use.input["mode"]
                mode_labels = {
                    "slideshow": "投影片影片",
                    "text_video": "文字動畫影片",
                    "tts_video": "TTS 語音影片",
                    "screen_record": f"螢幕錄影（{tool_use.input.get('duration',10)} 秒）",
                }
                await update.message.reply_text(f"🎬 生成中：{mode_labels.get(mode, mode)}，請稍候...")
                out_path = await loop.run_in_executor(None, lambda: execute_video_gen(
                    mode,
                    tool_use.input.get("output", ""),
                    text=tool_use.input.get("text", ""),
                    images=tool_use.input.get("images", []),
                    image=tool_use.input.get("image", ""),
                    duration=tool_use.input.get("duration", 5),
                    fps=tool_use.input.get("fps", 24),
                    voice=tool_use.input.get("voice", "zh-CN-YunxiNeural"),
                    bg_color=tool_use.input.get("bg_color", [30, 30, 40]),
                    font_color=tool_use.input.get("font_color", [255, 255, 255]),
                    font_size=tool_use.input.get("font_size", 60),
                    subtitle=tool_use.input.get("subtitle", True),
                ))
                if out_path and out_path.startswith("✅") and Path(out_path.split("：")[-1].strip()).exists():
                    video_file = out_path.split("：")[-1].strip()
                    await update.message.reply_chat_action("upload_video")
                    with open(video_file, "rb") as f:
                        await update.message.reply_video(
                            video=f,
                            caption=f"🎬 {mode_labels.get(mode, mode)}",
                            supports_streaming=True,
                        )
                    reply = out_path
                else:
                    reply = out_path
                save_message(chat_id, "assistant", reply)
                await _fr(reply)
                return

            elif tool_use.name == "sysres_chart":
                import asyncio
                loop = asyncio.get_running_loop()
                duration = tool_use.input.get("duration", 10)
                await update.message.reply_text(f"⏳ 監控 {duration} 秒中...")
                out_path = await loop.run_in_executor(None, execute_sysres_chart, duration)
                if out_path and Path(out_path).exists():
                    with open(out_path, "rb") as f:
                        await update.message.reply_photo(photo=f, caption="系統資源使用率")
                    reply = f"資源圖表已生成"
                else:
                    reply = out_path
                save_message(chat_id, "assistant", reply)
                await _fr(reply)
                return

            elif tool_use.name == "chart":
                import asyncio
                loop = asyncio.get_running_loop()
                out_path = await loop.run_in_executor(None, execute_chart,
                    tool_use.input["type"],
                    tool_use.input["data"],
                    tool_use.input.get("title",""),
                    tool_use.input.get("output",""))
                if out_path and Path(out_path).exists():
                    with open(out_path, "rb") as f:
                        await update.message.reply_photo(photo=f, caption=tool_use.input.get("title","圖表"))
                    reply = f"圖表已生成：{out_path}"
                else:
                    reply = out_path
                save_message(chat_id, "assistant", reply)
                await _fr(reply)
                return

            elif tool_use.name == "screen_stream":
                duration = tool_use.input.get("duration", 10)
                interval = tool_use.input.get("interval", 2)
                import asyncio
                loop = asyncio.get_running_loop()
                await update.message.reply_text(f"📹 開始串流 {duration} 秒...")
                screenshots = await loop.run_in_executor(None, execute_screen_stream, duration, interval)
                for i, img_bytes in enumerate(screenshots, 1):
                    await update.message.reply_photo(photo=img_bytes, caption=f"📸 {i}/{len(screenshots)}")
                reply = f"串流完成，共傳送 {len(screenshots)} 張截圖"
                save_message(chat_id, "assistant", reply)
                await _fr(reply)
                return

            elif tool_use.name in simple_tools:
                import asyncio as _asyncio
                _loop = _asyncio.get_running_loop()

                # 多步驟工具循環（最多10步），支援「打開記事本→輸入文字」等連續動作
                _messages = list(history)
                _cur_resp = response
                _last_result = ""

                for _step in range(10):
                    # 執行工具（lambdas 直接引用外層 tool_use 變數，更新 tool_use 即可）
                    _last_result = await _loop.run_in_executor(None, simple_tools[tool_use.name])

                    # open_app 動作後等視窗開啟
                    if tool_use.name == "desktop_control" and tool_use.input.get("action") == "open_app":
                        await _asyncio.sleep(1.0)

                    # 把結果傳回 Claude
                    _messages = _messages + [
                        {"role": "assistant", "content": _cur_resp.content},
                        {"role": "user", "content": _build_tool_results(_cur_resp.content, tool_use.id, _last_result)}
                    ]
                    _next_resp = client.messages.create(
                        model="claude-sonnet-4-6",
                        max_tokens=512,
                        system=system,
                        tools=TOOLS,
                        messages=_messages
                    )

                    # Claude 還要繼續呼叫 simple_tools？繼續循環
                    if _next_resp.stop_reason == "tool_use":
                        _next_tool = next((b for b in _next_resp.content if b.type == "tool_use"), None)
                        if _next_tool and _next_tool.name in simple_tools:
                            tool_use = _next_tool   # 更新 tool_use，lambdas 自動捕捉新值
                            _cur_resp = _next_resp
                            continue

                    # 結束循環
                    _cur_resp = _next_resp
                    break

                text_blocks = [b.text for b in _cur_resp.content if hasattr(b, "text")]
                reply = text_blocks[0] if text_blocks else _last_result
                save_message(chat_id, "assistant", reply)
                await _fr(reply)
                return

            elif tool_use.name == "tg_auto_reply":
                action = tool_use.input.get("action", "start")
                duration = tool_use.input.get("duration_minutes", 30)
                st = tool_use.input.get("stop_time", "")
                result_text = execute_tg_auto_reply(action, duration, st)
                save_message(chat_id, "assistant", result_text)
                await _fr(result_text)
                return

            elif tool_use.name == "screen_vision":
                question = tool_use.input.get("question", "請描述這個畫面上有什麼，以及目前電腦在做什麼事。")
                import asyncio
                loop = asyncio.get_running_loop()
                analysis, img_bytes = await loop.run_in_executor(None, execute_screen_vision, question)
                response = client.messages.create(
                    model="claude-sonnet-4-6",
                    max_tokens=1024,
                    system=system,
                    tools=TOOLS,
                    messages=history + [
                        {"role": "assistant", "content": response.content},
                        {"role": "user", "content": _build_tool_results(response.content, tool_use.id, analysis)}
                    ]
                )
                await update.message.reply_photo(photo=img_bytes, caption="📸 螢幕截圖")
                text_blocks = [b.text for b in response.content if hasattr(b, "text")]
                reply = text_blocks[0] if text_blocks else analysis
                save_message(chat_id, "assistant", reply)
                await _fr(reply)
                return

            elif tool_use.name == "find_image_on_screen":
                import asyncio
                loop = asyncio.get_running_loop()
                result = await loop.run_in_executor(
                    None, execute_find_image,
                    tool_use.input["template_path"],
                    tool_use.input.get("confidence", 0.8)
                )
                response = client.messages.create(
                    model="claude-sonnet-4-6",
                    max_tokens=512,
                    system=system,
                    tools=TOOLS,
                    messages=history + [
                        {"role": "assistant", "content": response.content},
                        {"role": "user", "content": _build_tool_results(response.content, tool_use.id, result)}
                    ]
                )
                text_blocks = [b.text for b in response.content if hasattr(b, "text")]
                reply = text_blocks[0] if text_blocks else result
                save_message(chat_id, "assistant", reply)
                await _fr(reply)
                return

            elif tool_use.name == "browser_control":
                inp = tool_use.input
                import asyncio
                loop = asyncio.get_running_loop()
                # 攔截 open/goto：改用 web_scrape 實際讀取內容（排除 YouTube 等影音網站）
                _url = inp.get("url", "")
                _skip_scrape = any(s in _url.lower() for s in ["youtube.com", "youtu.be", "netflix.com", "spotify.com", "twitch.tv"])
                if inp.get("action") in ("open", "goto") and _url and not _skip_scrape:
                    result = await loop.run_in_executor(
                        None, lambda: execute_web_scrape("scrape", url=_url)
                    )
                else:
                    result = await loop.run_in_executor(
                        None, lambda: execute_browser_control(
                            action=inp["action"],
                            url=inp.get("url", ""),
                            selector=inp.get("selector", ""),
                            text=inp.get("text", "")
                        )
                    )
                if result.startswith("__BROWSER_SCREENSHOT__:"):
                    img_bytes = bytes.fromhex(result.split(":", 1)[1])
                    await update.message.reply_photo(photo=img_bytes, caption="🌐 瀏覽器截圖")
                    tool_result_text = "已截圖並傳送"
                else:
                    tool_result_text = result
                response = client.messages.create(
                    model="claude-sonnet-4-6",
                    max_tokens=512,
                    system=system,
                    tools=TOOLS,
                    messages=history + [
                        {"role": "assistant", "content": response.content},
                        {"role": "user", "content": _build_tool_results(response.content, tool_use.id, tool_result_text)}
                    ]
                )
                text_blocks = [b.text for b in response.content if hasattr(b, "text")]
                reply = text_blocks[0] if text_blocks else tool_result_text
                save_message(chat_id, "assistant", reply)
                await _fr(reply)
                return

            elif tool_use.name == "long_term_memory":
                inp = tool_use.input
                action = inp["action"]
                if action == "save":
                    save_long_term_memory(chat_id, inp.get("content", ""))
                    tool_result = f"已記住：{inp.get('content', '')}"
                elif action == "list":
                    mems = load_long_term_memory(chat_id)
                    if mems:
                        tool_result = "長期記憶列表：\n" + "\n".join(f"[{m['id']}] {m['content']}" for m in mems)
                    else:
                        tool_result = "目前沒有長期記憶。"
                elif action == "delete":
                    delete_long_term_memory(inp.get("memory_id", 0))
                    tool_result = f"已刪除記憶 ID {inp.get('memory_id')}"
                else:
                    tool_result = "未知動作"
                response = client.messages.create(
                    model="claude-sonnet-4-6",
                    max_tokens=512,
                    system=system,
                    tools=TOOLS,
                    messages=history + [
                        {"role": "assistant", "content": response.content},
                        {"role": "user", "content": _build_tool_results(response.content, tool_use.id, tool_result)}
                    ]
                )
                text_blocks = [b.text for b in response.content if hasattr(b, "text")]
                reply = text_blocks[0] if text_blocks else tool_result
                save_message(chat_id, "assistant", reply)
                await _fr(reply)
                return

            elif tool_use.name == "get_stock":
                import asyncio
                loop = asyncio.get_running_loop()
                # Claude 可能同時查詢多支股票，全部執行並回傳
                all_stock_uses = [b for b in response.content
                                  if getattr(b, 'type', None) == "tool_use" and b.name == "get_stock"]
                stock_results = await asyncio.gather(*[
                    loop.run_in_executor(None, fetch_stock, tu.input["symbol"], tu.input.get("period", "1mo"))
                    for tu in all_stock_uses
                ])
                tool_results_list = [
                    {"type": "tool_result", "tool_use_id": tu.id, "content": res}
                    for tu, res in zip(all_stock_uses, stock_results)
                ]
                # 補齊非 get_stock 的其他 tool_use（如有）
                for b in response.content:
                    if getattr(b, 'type', None) == "tool_use" and b.name != "get_stock":
                        tool_results_list.append({"type": "tool_result", "tool_use_id": b.id, "content": "（跳過）"})
                response = client.messages.create(
                    model="claude-sonnet-4-6",
                    max_tokens=1024,
                    system=system,
                    tools=TOOLS,
                    messages=history + [
                        {"role": "assistant", "content": response.content},
                        {"role": "user", "content": tool_results_list}
                    ]
                )

            elif tool_use.name == "get_weather":
                tool_result = fetch_weather(tool_use.input["city"])
                response = client.messages.create(
                    model="claude-sonnet-4-6",
                    max_tokens=1024,
                    system=system,
                    tools=TOOLS,
                    messages=history + [
                        {"role": "assistant", "content": response.content},
                        {"role": "user", "content": _build_tool_results(response.content, tool_use.id, tool_result)}
                    ]
                )

            elif tool_use.name == "get_crypto":
                import asyncio
                loop = asyncio.get_running_loop()
                tool_result = await loop.run_in_executor(None, fetch_crypto,
                    tool_use.input["coin"], tool_use.input.get("vs_currency", "usd"))
                response = client.messages.create(
                    model="claude-sonnet-4-6", max_tokens=1024, system=system, tools=TOOLS,
                    messages=history + [
                        {"role": "assistant", "content": response.content},
                        {"role": "user", "content": _build_tool_results(response.content, tool_use.id, tool_result)}
                    ]
                )

            elif tool_use.name == "get_forex":
                import asyncio
                loop = asyncio.get_running_loop()
                tool_result = await loop.run_in_executor(None, fetch_forex, tool_use.input["pair"])
                response = client.messages.create(
                    model="claude-sonnet-4-6", max_tokens=1024, system=system, tools=TOOLS,
                    messages=history + [
                        {"role": "assistant", "content": response.content},
                        {"role": "user", "content": _build_tool_results(response.content, tool_use.id, tool_result)}
                    ]
                )

            elif tool_use.name == "get_finance_news":
                import asyncio
                loop = asyncio.get_running_loop()
                tool_result = await loop.run_in_executor(None, fetch_finance_news,
                    tool_use.input.get("source", "all"), tool_use.input.get("count", 5))
                response = client.messages.create(
                    model="claude-sonnet-4-6", max_tokens=1024, system=system, tools=TOOLS,
                    messages=history + [
                        {"role": "assistant", "content": response.content},
                        {"role": "user", "content": _build_tool_results(response.content, tool_use.id, tool_result)}
                    ]
                )

            elif tool_use.name == "china_search":
                import asyncio
                loop = asyncio.get_running_loop()
                tool_result = await loop.run_in_executor(None, fetch_china_search,
                    tool_use.input["query"],
                    tool_use.input.get("category", "其他"),
                    tool_use.input.get("count", 6))
                response = client.messages.create(
                    model="claude-sonnet-4-6", max_tokens=1500, system=system, tools=TOOLS,
                    messages=history + [
                        {"role": "assistant", "content": response.content},
                        {"role": "user", "content": _build_tool_results(response.content, tool_use.id, tool_result)}
                    ]
                )

            elif tool_use.name == "get_ashare":
                import asyncio
                loop = asyncio.get_running_loop()
                tool_result = await loop.run_in_executor(None, fetch_ashare,
                    tool_use.input["code"], tool_use.input.get("period", "1mo"))
                response = client.messages.create(
                    model="claude-sonnet-4-6", max_tokens=1024, system=system, tools=TOOLS,
                    messages=history + [
                        {"role": "assistant", "content": response.content},
                        {"role": "user", "content": _build_tool_results(response.content, tool_use.id, tool_result)}
                    ]
                )

            elif tool_use.name == "get_cn_news":
                import asyncio
                loop = asyncio.get_running_loop()
                tool_result = await loop.run_in_executor(None, fetch_cn_news,
                    tool_use.input.get("source", "all"), tool_use.input.get("count", 5))
                response = client.messages.create(
                    model="claude-sonnet-4-6", max_tokens=1024, system=system, tools=TOOLS,
                    messages=history + [
                        {"role": "assistant", "content": response.content},
                        {"role": "user", "content": _build_tool_results(response.content, tool_use.id, tool_result)}
                    ]
                )

            elif tool_use.name == "get_institutional":
                import asyncio
                loop = asyncio.get_running_loop()
                tool_result = await loop.run_in_executor(None, fetch_institutional,
                    tool_use.input.get("symbol", ""), tool_use.input.get("date", ""))
                response = client.messages.create(
                    model="claude-sonnet-4-6", max_tokens=1024, system=system, tools=TOOLS,
                    messages=history + [{"role": "assistant", "content": response.content},
                    {"role": "user", "content": _build_tool_results(response.content, tool_use.id, tool_result)}])

            elif tool_use.name == "get_sector":
                import asyncio
                loop = asyncio.get_running_loop()
                tool_result = await loop.run_in_executor(None, fetch_sector,
                    tool_use.input.get("market", "us"))
                response = client.messages.create(
                    model="claude-sonnet-4-6", max_tokens=1024, system=system, tools=TOOLS,
                    messages=history + [{"role": "assistant", "content": response.content},
                    {"role": "user", "content": _build_tool_results(response.content, tool_use.id, tool_result)}])

            elif tool_use.name == "get_commodity":
                import asyncio
                loop = asyncio.get_running_loop()
                tool_result = await loop.run_in_executor(None, fetch_commodity,
                    tool_use.input.get("items", ["all"]))
                response = client.messages.create(
                    model="claude-sonnet-4-6", max_tokens=1024, system=system, tools=TOOLS,
                    messages=history + [{"role": "assistant", "content": response.content},
                    {"role": "user", "content": _build_tool_results(response.content, tool_use.id, tool_result)}])

            elif tool_use.name == "get_bond_yield":
                import asyncio
                loop = asyncio.get_running_loop()
                tool_result = await loop.run_in_executor(None, fetch_bond_yield)
                response = client.messages.create(
                    model="claude-sonnet-4-6", max_tokens=1024, system=system, tools=TOOLS,
                    messages=history + [{"role": "assistant", "content": response.content},
                    {"role": "user", "content": _build_tool_results(response.content, tool_use.id, tool_result)}])

            elif tool_use.name == "get_dividend_calendar":
                import asyncio
                loop = asyncio.get_running_loop()
                tool_result = await loop.run_in_executor(None, fetch_dividend_calendar,
                    tool_use.input["symbol"])
                response = client.messages.create(
                    model="claude-sonnet-4-6", max_tokens=1024, system=system, tools=TOOLS,
                    messages=history + [{"role": "assistant", "content": response.content},
                    {"role": "user", "content": _build_tool_results(response.content, tool_use.id, tool_result)}])

            elif tool_use.name == "stock_screener":
                import asyncio
                loop = asyncio.get_running_loop()
                tool_result = await loop.run_in_executor(None, fetch_stock_screener,
                    tool_use.input["criteria"], tool_use.input.get("market", "us"))
                response = client.messages.create(
                    model="claude-sonnet-4-6", max_tokens=1500, system=system, tools=TOOLS,
                    messages=history + [{"role": "assistant", "content": response.content},
                    {"role": "user", "content": _build_tool_results(response.content, tool_use.id, tool_result)}])

            elif tool_use.name == "get_margin_trading":
                import asyncio
                loop = asyncio.get_running_loop()
                tool_result = await loop.run_in_executor(None, fetch_margin_trading,
                    tool_use.input["symbol"], tool_use.input.get("date", ""))
                response = client.messages.create(
                    model="claude-sonnet-4-6", max_tokens=1024, system=system, tools=TOOLS,
                    messages=history + [{"role": "assistant", "content": response.content},
                    {"role": "user", "content": _build_tool_results(response.content, tool_use.id, tool_result)}])

            elif tool_use.name == "get_options":
                import asyncio
                loop = asyncio.get_running_loop()
                tool_result = await loop.run_in_executor(None, fetch_options,
                    tool_use.input["symbol"], tool_use.input.get("expiry", ""))
                response = client.messages.create(
                    model="claude-sonnet-4-6", max_tokens=1024, system=system, tools=TOOLS,
                    messages=history + [{"role": "assistant", "content": response.content},
                    {"role": "user", "content": _build_tool_results(response.content, tool_use.id, tool_result)}])

            elif tool_use.name == "get_futures":
                import asyncio
                loop = asyncio.get_running_loop()
                tool_result = await loop.run_in_executor(None, fetch_futures,
                    tool_use.input.get("items", ["all"]))
                response = client.messages.create(
                    model="claude-sonnet-4-6", max_tokens=1024, system=system, tools=TOOLS,
                    messages=history + [{"role": "assistant", "content": response.content},
                    {"role": "user", "content": _build_tool_results(response.content, tool_use.id, tool_result)}])

            elif tool_use.name == "get_ipo":
                import asyncio
                loop = asyncio.get_running_loop()
                tool_result = await loop.run_in_executor(None, fetch_ipo,
                    tool_use.input.get("count", 10))
                response = client.messages.create(
                    model="claude-sonnet-4-6", max_tokens=1024, system=system, tools=TOOLS,
                    messages=history + [{"role": "assistant", "content": response.content},
                    {"role": "user", "content": _build_tool_results(response.content, tool_use.id, tool_result)}])

            elif tool_use.name == "backtest":
                import asyncio
                loop = asyncio.get_running_loop()
                tool_result = await loop.run_in_executor(None, fetch_backtest,
                    tool_use.input["symbol"],
                    tool_use.input.get("strategy", "ma_cross"),
                    tool_use.input.get("period", "2y"))
                response = client.messages.create(
                    model="claude-sonnet-4-6", max_tokens=1024, system=system, tools=TOOLS,
                    messages=history + [{"role": "assistant", "content": response.content},
                    {"role": "user", "content": _build_tool_results(response.content, tool_use.id, tool_result)}])

            elif tool_use.name == "get_global_market":
                import asyncio; loop = asyncio.get_running_loop()
                tool_result = await loop.run_in_executor(None, fetch_global_market)
                response = client.messages.create(model="claude-sonnet-4-6", max_tokens=1024, system=system, tools=TOOLS,
                    messages=history + [{"role": "assistant", "content": response.content},
                    {"role": "user", "content": _build_tool_results(response.content, tool_use.id, tool_result)}])

            elif tool_use.name == "get_economic_calendar":
                import asyncio; loop = asyncio.get_running_loop()
                tool_result = await loop.run_in_executor(None, fetch_economic_calendar,
                    tool_use.input.get("count", 10))
                response = client.messages.create(model="claude-sonnet-4-6", max_tokens=1024, system=system, tools=TOOLS,
                    messages=history + [{"role": "assistant", "content": response.content},
                    {"role": "user", "content": _build_tool_results(response.content, tool_use.id, tool_result)}])

            elif tool_use.name == "get_earnings_calendar":
                import asyncio; loop = asyncio.get_running_loop()
                tool_result = await loop.run_in_executor(None, fetch_earnings_calendar,
                    tool_use.input.get("days", 7))
                response = client.messages.create(model="claude-sonnet-4-6", max_tokens=1024, system=system, tools=TOOLS,
                    messages=history + [{"role": "assistant", "content": response.content},
                    {"role": "user", "content": _build_tool_results(response.content, tool_use.id, tool_result)}])

            elif tool_use.name == "get_analyst_ratings":
                import asyncio; loop = asyncio.get_running_loop()
                tool_result = await loop.run_in_executor(None, fetch_analyst_ratings,
                    tool_use.input["symbol"])
                response = client.messages.create(model="claude-sonnet-4-6", max_tokens=1024, system=system, tools=TOOLS,
                    messages=history + [{"role": "assistant", "content": response.content},
                    {"role": "user", "content": _build_tool_results(response.content, tool_use.id, tool_result)}])

            elif tool_use.name == "get_short_interest":
                import asyncio; loop = asyncio.get_running_loop()
                tool_result = await loop.run_in_executor(None, fetch_short_interest,
                    tool_use.input["symbol"])
                response = client.messages.create(model="claude-sonnet-4-6", max_tokens=1024, system=system, tools=TOOLS,
                    messages=history + [{"role": "assistant", "content": response.content},
                    {"role": "user", "content": _build_tool_results(response.content, tool_use.id, tool_result)}])

            elif tool_use.name == "get_correlation":
                import asyncio; loop = asyncio.get_running_loop()
                tool_result = await loop.run_in_executor(None, fetch_correlation,
                    tool_use.input["symbols"], tool_use.input.get("period", "1y"))
                response = client.messages.create(model="claude-sonnet-4-6", max_tokens=1024, system=system, tools=TOOLS,
                    messages=history + [{"role": "assistant", "content": response.content},
                    {"role": "user", "content": _build_tool_results(response.content, tool_use.id, tool_result)}])

            elif tool_use.name == "get_risk_metrics":
                import asyncio; loop = asyncio.get_running_loop()
                tool_result = await loop.run_in_executor(None, fetch_risk_metrics,
                    tool_use.input["symbol"], tool_use.input.get("period", "1y"))
                response = client.messages.create(model="claude-sonnet-4-6", max_tokens=1024, system=system, tools=TOOLS,
                    messages=history + [{"role": "assistant", "content": response.content},
                    {"role": "user", "content": _build_tool_results(response.content, tool_use.id, tool_result)}])

            elif tool_use.name == "get_money_flow":
                import asyncio; loop = asyncio.get_running_loop()
                tool_result = await loop.run_in_executor(None, fetch_money_flow,
                    tool_use.input["symbol"])
                response = client.messages.create(model="claude-sonnet-4-6", max_tokens=1024, system=system, tools=TOOLS,
                    messages=history + [{"role": "assistant", "content": response.content},
                    {"role": "user", "content": _build_tool_results(response.content, tool_use.id, tool_result)}])

            elif tool_use.name == "get_concept_stocks":
                import asyncio; loop = asyncio.get_running_loop()
                tool_result = await loop.run_in_executor(None, fetch_concept_stocks,
                    tool_use.input["theme"])
                response = client.messages.create(model="claude-sonnet-4-6", max_tokens=1500, system=system, tools=TOOLS,
                    messages=history + [{"role": "assistant", "content": response.content},
                    {"role": "user", "content": _build_tool_results(response.content, tool_use.id, tool_result)}])

            elif tool_use.name == "get_crypto_depth":
                import asyncio; loop = asyncio.get_running_loop()
                tool_result = await loop.run_in_executor(None, fetch_crypto_depth,
                    tool_use.input.get("coin", "bitcoin"))
                response = client.messages.create(model="claude-sonnet-4-6", max_tokens=1024, system=system, tools=TOOLS,
                    messages=history + [{"role": "assistant", "content": response.content},
                    {"role": "user", "content": _build_tool_results(response.content, tool_use.id, tool_result)}])

            elif tool_use.name == "drip_calculator":
                import asyncio; loop = asyncio.get_running_loop()
                tool_result = await loop.run_in_executor(None, fetch_drip_calculator,
                    tool_use.input["symbol"], tool_use.input["shares"],
                    tool_use.input.get("years", 10), tool_use.input.get("monthly_invest", 0))
                response = client.messages.create(model="claude-sonnet-4-6", max_tokens=1024, system=system, tools=TOOLS,
                    messages=history + [{"role": "assistant", "content": response.content},
                    {"role": "user", "content": _build_tool_results(response.content, tool_use.id, tool_result)}])

            elif tool_use.name == "get_forex_chart":
                import asyncio; loop = asyncio.get_running_loop()
                tool_result = await loop.run_in_executor(None, fetch_forex_chart,
                    tool_use.input["pair"], tool_use.input.get("period", "3mo"))
                response = client.messages.create(model="claude-sonnet-4-6", max_tokens=1024, system=system, tools=TOOLS,
                    messages=history + [{"role": "assistant", "content": response.content},
                    {"role": "user", "content": _build_tool_results(response.content, tool_use.id, tool_result)}])

            elif tool_use.name == "get_warrant":
                import asyncio; loop = asyncio.get_running_loop()
                tool_result = await loop.run_in_executor(None, fetch_warrant,
                    tool_use.input["underlying"])
                response = client.messages.create(model="claude-sonnet-4-6", max_tokens=1024, system=system, tools=TOOLS,
                    messages=history + [{"role": "assistant", "content": response.content},
                    {"role": "user", "content": _build_tool_results(response.content, tool_use.id, tool_result)}])

            elif tool_use.name == "get_portfolio_risk":
                import asyncio; loop = asyncio.get_running_loop()
                tool_result = await loop.run_in_executor(None, fetch_portfolio_risk,
                    tool_use.input["holdings"], tool_use.input.get("period", "1y"))
                response = client.messages.create(model="claude-sonnet-4-6", max_tokens=1500, system=system, tools=TOOLS,
                    messages=history + [{"role": "assistant", "content": response.content},
                    {"role": "user", "content": _build_tool_results(response.content, tool_use.id, tool_result)}])

            elif tool_use.name == "retirement_calculator":
                import asyncio; loop = asyncio.get_running_loop()
                tool_result = await loop.run_in_executor(None, fetch_retirement_calculator,
                    tool_use.input["current_age"], tool_use.input["current_savings"],
                    tool_use.input["monthly_save"], tool_use.input.get("retire_age", 65),
                    tool_use.input.get("annual_return", 6.0), tool_use.input.get("monthly_expense", 50000))
                response = client.messages.create(model="claude-sonnet-4-6", max_tokens=1024, system=system, tools=TOOLS,
                    messages=history + [{"role": "assistant", "content": response.content},
                    {"role": "user", "content": _build_tool_results(response.content, tool_use.id, tool_result)}])

            elif tool_use.name == "loan_calculator":
                import asyncio; loop = asyncio.get_running_loop()
                tool_result = await loop.run_in_executor(None, fetch_loan_calculator,
                    tool_use.input["principal"], tool_use.input["annual_rate"],
                    tool_use.input["years"], tool_use.input.get("loan_type", "等額本息"))
                response = client.messages.create(model="claude-sonnet-4-6", max_tokens=1024, system=system, tools=TOOLS,
                    messages=history + [{"role": "assistant", "content": response.content},
                    {"role": "user", "content": _build_tool_results(response.content, tool_use.id, tool_result)}])

            elif tool_use.name == "compound_calculator":
                import asyncio; loop = asyncio.get_running_loop()
                tool_result = await loop.run_in_executor(None, fetch_compound_calculator,
                    tool_use.input["principal"], tool_use.input["annual_rate"],
                    tool_use.input["years"], tool_use.input.get("monthly_add", 0),
                    tool_use.input.get("compound_freq", 12))
                response = client.messages.create(model="claude-sonnet-4-6", max_tokens=1024, system=system, tools=TOOLS,
                    messages=history + [{"role": "assistant", "content": response.content},
                    {"role": "user", "content": _build_tool_results(response.content, tool_use.id, tool_result)}])

            elif tool_use.name == "asset_allocation":
                import asyncio; loop = asyncio.get_running_loop()
                tool_result = await loop.run_in_executor(None, fetch_asset_allocation,
                    tool_use.input["age"], tool_use.input["risk_level"],
                    tool_use.input.get("goal", "退休"), tool_use.input.get("investment_horizon"))
                response = client.messages.create(model="claude-sonnet-4-6", max_tokens=1024, system=system, tools=TOOLS,
                    messages=history + [{"role": "assistant", "content": response.content},
                    {"role": "user", "content": _build_tool_results(response.content, tool_use.id, tool_result)}])

            elif tool_use.name == "tw_tax_calculator":
                import asyncio; loop = asyncio.get_running_loop()
                tool_result = await loop.run_in_executor(None, fetch_tw_tax_calculator,
                    tool_use.input["dividend_income"], tool_use.input.get("other_income", 0),
                    tool_use.input.get("tax_bracket"), tool_use.input.get("sell_amount", 0))
                response = client.messages.create(model="claude-sonnet-4-6", max_tokens=1024, system=system, tools=TOOLS,
                    messages=history + [{"role": "assistant", "content": response.content},
                    {"role": "user", "content": _build_tool_results(response.content, tool_use.id, tool_result)}])

            elif tool_use.name == "currency_converter":
                import asyncio; loop = asyncio.get_running_loop()
                tool_result = await loop.run_in_executor(None, fetch_currency_converter,
                    tool_use.input["amount"], tool_use.input["from_currency"],
                    tool_use.input["to_currency"])
                response = client.messages.create(model="claude-sonnet-4-6", max_tokens=512, system=system, tools=TOOLS,
                    messages=history + [{"role": "assistant", "content": response.content},
                    {"role": "user", "content": _build_tool_results(response.content, tool_use.id, tool_result)}])

            elif tool_use.name == "get_fund":
                import asyncio; loop = asyncio.get_running_loop()
                tool_result = await loop.run_in_executor(None, fetch_fund,
                    tool_use.input["symbol"])
                response = client.messages.create(model="claude-sonnet-4-6", max_tokens=1024, system=system, tools=TOOLS,
                    messages=history + [{"role": "assistant", "content": response.content},
                    {"role": "user", "content": _build_tool_results(response.content, tool_use.id, tool_result)}])

            elif tool_use.name == "get_reits":
                import asyncio; loop = asyncio.get_running_loop()
                tool_result = await loop.run_in_executor(None, fetch_reits,
                    tool_use.input["symbol"])
                response = client.messages.create(model="claude-sonnet-4-6", max_tokens=1024, system=system, tools=TOOLS,
                    messages=history + [{"role": "assistant", "content": response.content},
                    {"role": "user", "content": _build_tool_results(response.content, tool_use.id, tool_result)}])

            elif tool_use.name == "inflation_adjusted":
                import asyncio; loop = asyncio.get_running_loop()
                tool_result = await loop.run_in_executor(None, fetch_inflation_adjusted,
                    tool_use.input["nominal_return"], tool_use.input["years"],
                    tool_use.input["amount"], tool_use.input.get("inflation_rate", 2.0))
                response = client.messages.create(model="claude-sonnet-4-6", max_tokens=1024, system=system, tools=TOOLS,
                    messages=history + [{"role": "assistant", "content": response.content},
                    {"role": "user", "content": _build_tool_results(response.content, tool_use.id, tool_result)}])

            elif tool_use.name == "defi_calculator":
                import asyncio; loop = asyncio.get_running_loop()
                tool_result = await loop.run_in_executor(None, fetch_defi_calculator,
                    tool_use.input["principal_usd"], tool_use.input["apy"],
                    tool_use.input["days"], tool_use.input.get("compound", True),
                    tool_use.input.get("protocol", ""))
                response = client.messages.create(model="claude-sonnet-4-6", max_tokens=1024, system=system, tools=TOOLS,
                    messages=history + [{"role": "assistant", "content": response.content},
                    {"role": "user", "content": _build_tool_results(response.content, tool_use.id, tool_result)}])

            elif tool_use.name == "gold_calculator":
                import asyncio; loop = asyncio.get_running_loop()
                tool_result = await loop.run_in_executor(None, fetch_gold_calculator,
                    tool_use.input["weight"], tool_use.input.get("unit", "公克"),
                    tool_use.input.get("currency", "TWD"))
                response = client.messages.create(model="claude-sonnet-4-6", max_tokens=512, system=system, tools=TOOLS,
                    messages=history + [{"role": "assistant", "content": response.content},
                    {"role": "user", "content": _build_tool_results(response.content, tool_use.id, tool_result)}])

            elif tool_use.name == "forex_deposit":
                import asyncio; loop = asyncio.get_running_loop()
                tool_result = await loop.run_in_executor(None, fetch_forex_deposit,
                    tool_use.input["amount_twd"], tool_use.input["currency"],
                    tool_use.input["annual_rate"], tool_use.input["months"],
                    tool_use.input.get("buy_rate"), tool_use.input.get("sell_rate"))
                response = client.messages.create(model="claude-sonnet-4-6", max_tokens=1024, system=system, tools=TOOLS,
                    messages=history + [{"role": "assistant", "content": response.content},
                    {"role": "user", "content": _build_tool_results(response.content, tool_use.id, tool_result)}])

            elif tool_use.name == "financial_health":
                import asyncio; loop = asyncio.get_running_loop()
                tool_result = await loop.run_in_executor(None, fetch_financial_health,
                    tool_use.input["monthly_income"], tool_use.input["monthly_expense"],
                    tool_use.input["total_assets"], tool_use.input["total_debt"],
                    tool_use.input.get("emergency_fund_months", 0),
                    tool_use.input.get("has_insurance", False),
                    tool_use.input.get("investment_ratio", 0))
                response = client.messages.create(model="claude-sonnet-4-6", max_tokens=1024, system=system, tools=TOOLS,
                    messages=history + [{"role": "assistant", "content": response.content},
                    {"role": "user", "content": _build_tool_results(response.content, tool_use.id, tool_result)}])

            elif tool_use.name == "deep_research":
                import asyncio; loop = asyncio.get_running_loop()
                tool_result = await loop.run_in_executor(None, fetch_deep_research,
                    tool_use.input["topic"], tool_use.input.get("lang", "zh-tw"),
                    tool_use.input.get("depth", 5))
                response = client.messages.create(model="claude-sonnet-4-6", max_tokens=2048, system=system, tools=TOOLS,
                    messages=history + [{"role": "assistant", "content": response.content},
                    {"role": "user", "content": _build_tool_results(response.content, tool_use.id, tool_result)}])

            elif tool_use.name == "fact_check":
                import asyncio; loop = asyncio.get_running_loop()
                tool_result = await loop.run_in_executor(None, fetch_fact_check,
                    tool_use.input["claim"], tool_use.input.get("lang", "zh-tw"))
                response = client.messages.create(model="claude-sonnet-4-6", max_tokens=1536, system=system, tools=TOOLS,
                    messages=history + [{"role": "assistant", "content": response.content},
                    {"role": "user", "content": _build_tool_results(response.content, tool_use.id, tool_result)}])

            elif tool_use.name == "timeline_events":
                import asyncio; loop = asyncio.get_running_loop()
                tool_result = await loop.run_in_executor(None, fetch_timeline_events,
                    tool_use.input["topic"], tool_use.input.get("lang", "zh-tw"))
                response = client.messages.create(model="claude-sonnet-4-6", max_tokens=1536, system=system, tools=TOOLS,
                    messages=history + [{"role": "assistant", "content": response.content},
                    {"role": "user", "content": _build_tool_results(response.content, tool_use.id, tool_result)}])

            elif tool_use.name == "sentiment_scan":
                import asyncio; loop = asyncio.get_running_loop()
                tool_result = await loop.run_in_executor(None, fetch_sentiment_scan,
                    tool_use.input["topic"], tool_use.input.get("lang", "zh-tw"))
                response = client.messages.create(model="claude-sonnet-4-6", max_tokens=1536, system=system, tools=TOOLS,
                    messages=history + [{"role": "assistant", "content": response.content},
                    {"role": "user", "content": _build_tool_results(response.content, tool_use.id, tool_result)}])

            elif tool_use.name == "compare_analysis":
                import asyncio; loop = asyncio.get_running_loop()
                tool_result = await loop.run_in_executor(None, fetch_compare_analysis,
                    tool_use.input["items"], tool_use.input.get("dimensions"),
                    tool_use.input.get("context", ""))
                response = client.messages.create(model="claude-sonnet-4-6", max_tokens=2048, system=system, tools=TOOLS,
                    messages=history + [{"role": "assistant", "content": response.content},
                    {"role": "user", "content": _build_tool_results(response.content, tool_use.id, tool_result)}])

            elif tool_use.name == "pros_cons_analysis":
                import asyncio; loop = asyncio.get_running_loop()
                tool_result = await loop.run_in_executor(None, fetch_pros_cons_analysis,
                    tool_use.input["subject"], tool_use.input.get("context", ""),
                    tool_use.input.get("lang", "zh-tw"))
                response = client.messages.create(model="claude-sonnet-4-6", max_tokens=1536, system=system, tools=TOOLS,
                    messages=history + [{"role": "assistant", "content": response.content},
                    {"role": "user", "content": _build_tool_results(response.content, tool_use.id, tool_result)}])

            elif tool_use.name == "research_report":
                import asyncio; loop = asyncio.get_running_loop()
                tool_result = await loop.run_in_executor(None, fetch_research_report,
                    tool_use.input["topic"], tool_use.input.get("purpose", "一般研究"),
                    tool_use.input.get("lang", "zh-tw"))
                response = client.messages.create(model="claude-sonnet-4-6", max_tokens=2048, system=system, tools=TOOLS,
                    messages=history + [{"role": "assistant", "content": response.content},
                    {"role": "user", "content": _build_tool_results(response.content, tool_use.id, tool_result)}])

            elif tool_use.name == "opinion_writer":
                import asyncio; loop = asyncio.get_running_loop()
                tool_result = await loop.run_in_executor(None, fetch_opinion_writer,
                    tool_use.input["topic"], tool_use.input.get("stance", "中立"),
                    tool_use.input.get("style", "正式"))
                response = client.messages.create(model="claude-sonnet-4-6", max_tokens=2048, system=system, tools=TOOLS,
                    messages=history + [{"role": "assistant", "content": response.content},
                    {"role": "user", "content": _build_tool_results(response.content, tool_use.id, tool_result)}])

            elif tool_use.name == "trend_forecast":
                import asyncio; loop = asyncio.get_running_loop()
                tool_result = await loop.run_in_executor(None, fetch_trend_forecast,
                    tool_use.input["topic"], tool_use.input.get("timeframe", "全部"),
                    tool_use.input.get("lang", "zh-tw"))
                response = client.messages.create(model="claude-sonnet-4-6", max_tokens=2048, system=system, tools=TOOLS,
                    messages=history + [{"role": "assistant", "content": response.content},
                    {"role": "user", "content": _build_tool_results(response.content, tool_use.id, tool_result)}])

            elif tool_use.name == "debate_simulator":
                import asyncio; loop = asyncio.get_running_loop()
                tool_result = await loop.run_in_executor(None, fetch_debate_simulator,
                    tool_use.input["motion"], tool_use.input.get("lang", "zh-tw"))
                response = client.messages.create(model="claude-sonnet-4-6", max_tokens=2048, system=system, tools=TOOLS,
                    messages=history + [{"role": "assistant", "content": response.content},
                    {"role": "user", "content": _build_tool_results(response.content, tool_use.id, tool_result)}])

            elif tool_use.name == "academic_search":
                import asyncio; loop = asyncio.get_running_loop()
                tool_result = await loop.run_in_executor(None, fetch_academic_search,
                    tool_use.input["query"], tool_use.input.get("field", ""),
                    tool_use.input.get("lang", "en"))
                response = client.messages.create(model="claude-sonnet-4-6", max_tokens=2048, system=system, tools=TOOLS,
                    messages=history + [{"role": "assistant", "content": response.content},
                    {"role": "user", "content": _build_tool_results(response.content, tool_use.id, tool_result)}])

            elif tool_use.name == "health_research":
                import asyncio; loop = asyncio.get_running_loop()
                tool_result = await loop.run_in_executor(None, fetch_health_research,
                    tool_use.input["topic"], tool_use.input.get("lang", "zh-tw"))
                response = client.messages.create(model="claude-sonnet-4-6", max_tokens=1536, system=system, tools=TOOLS,
                    messages=history + [{"role": "assistant", "content": response.content},
                    {"role": "user", "content": _build_tool_results(response.content, tool_use.id, tool_result)}])

            elif tool_use.name == "law_research":
                import asyncio; loop = asyncio.get_running_loop()
                tool_result = await loop.run_in_executor(None, fetch_law_research,
                    tool_use.input["topic"], tool_use.input.get("jurisdiction", "台灣"),
                    tool_use.input.get("lang", "zh-tw"))
                response = client.messages.create(model="claude-sonnet-4-6", max_tokens=1536, system=system, tools=TOOLS,
                    messages=history + [{"role": "assistant", "content": response.content},
                    {"role": "user", "content": _build_tool_results(response.content, tool_use.id, tool_result)}])

            elif tool_use.name == "person_research":
                import asyncio; loop = asyncio.get_running_loop()
                tool_result = await loop.run_in_executor(None, fetch_person_research,
                    tool_use.input["name"], tool_use.input.get("context", ""),
                    tool_use.input.get("lang", "zh-tw"))
                response = client.messages.create(model="claude-sonnet-4-6", max_tokens=2048, system=system, tools=TOOLS,
                    messages=history + [{"role": "assistant", "content": response.content},
                    {"role": "user", "content": _build_tool_results(response.content, tool_use.id, tool_result)}])

            elif tool_use.name == "company_research":
                import asyncio; loop = asyncio.get_running_loop()
                tool_result = await loop.run_in_executor(None, fetch_company_research,
                    tool_use.input["company"], tool_use.input.get("lang", "zh-tw"))
                response = client.messages.create(model="claude-sonnet-4-6", max_tokens=2048, system=system, tools=TOOLS,
                    messages=history + [{"role": "assistant", "content": response.content},
                    {"role": "user", "content": _build_tool_results(response.content, tool_use.id, tool_result)}])

            elif tool_use.name == "product_review":
                import asyncio; loop = asyncio.get_running_loop()
                tool_result = await loop.run_in_executor(None, fetch_product_review,
                    tool_use.input["product"], tool_use.input.get("category", ""),
                    tool_use.input.get("lang", "zh-tw"))
                response = client.messages.create(model="claude-sonnet-4-6", max_tokens=1536, system=system, tools=TOOLS,
                    messages=history + [{"role": "assistant", "content": response.content},
                    {"role": "user", "content": _build_tool_results(response.content, tool_use.id, tool_result)}])

            elif tool_use.name == "travel_research":
                import asyncio; loop = asyncio.get_running_loop()
                tool_result = await loop.run_in_executor(None, fetch_travel_research,
                    tool_use.input["destination"], tool_use.input.get("days"),
                    tool_use.input.get("style", ""), tool_use.input.get("lang", "zh-tw"))
                response = client.messages.create(model="claude-sonnet-4-6", max_tokens=2048, system=system, tools=TOOLS,
                    messages=history + [{"role": "assistant", "content": response.content},
                    {"role": "user", "content": _build_tool_results(response.content, tool_use.id, tool_result)}])

            elif tool_use.name == "job_market":
                import asyncio; loop = asyncio.get_running_loop()
                tool_result = await loop.run_in_executor(None, fetch_job_market,
                    tool_use.input["job_title"], tool_use.input.get("location", "台灣"),
                    tool_use.input.get("lang", "zh-tw"))
                response = client.messages.create(model="claude-sonnet-4-6", max_tokens=1536, system=system, tools=TOOLS,
                    messages=history + [{"role": "assistant", "content": response.content},
                    {"role": "user", "content": _build_tool_results(response.content, tool_use.id, tool_result)}])

            elif tool_use.name == "impact_analysis":
                import asyncio; loop = asyncio.get_running_loop()
                tool_result = await loop.run_in_executor(None, fetch_impact_analysis,
                    tool_use.input["event"], tool_use.input.get("scope"),
                    tool_use.input.get("lang", "zh-tw"))
                response = client.messages.create(model="claude-sonnet-4-6", max_tokens=2048, system=system, tools=TOOLS,
                    messages=history + [{"role": "assistant", "content": response.content},
                    {"role": "user", "content": _build_tool_results(response.content, tool_use.id, tool_result)}])

            elif tool_use.name == "scenario_planning":
                import asyncio; loop = asyncio.get_running_loop()
                tool_result = await loop.run_in_executor(None, fetch_scenario_planning,
                    tool_use.input["topic"], tool_use.input.get("horizon", ""),
                    tool_use.input.get("lang", "zh-tw"))
                response = client.messages.create(model="claude-sonnet-4-6", max_tokens=2048, system=system, tools=TOOLS,
                    messages=history + [{"role": "assistant", "content": response.content},
                    {"role": "user", "content": _build_tool_results(response.content, tool_use.id, tool_result)}])

            elif tool_use.name == "decision_helper":
                import asyncio; loop = asyncio.get_running_loop()
                tool_result = await loop.run_in_executor(None, fetch_decision_helper,
                    tool_use.input["question"], tool_use.input.get("options"),
                    tool_use.input.get("criteria"))
                response = client.messages.create(model="claude-sonnet-4-6", max_tokens=2048, system=system, tools=TOOLS,
                    messages=history + [{"role": "assistant", "content": response.content},
                    {"role": "user", "content": _build_tool_results(response.content, tool_use.id, tool_result)}])

            elif tool_use.name == "devil_advocate":
                import asyncio; loop = asyncio.get_running_loop()
                tool_result = await loop.run_in_executor(None, fetch_devil_advocate,
                    tool_use.input["position"], tool_use.input.get("lang", "zh-tw"))
                response = client.messages.create(model="claude-sonnet-4-6", max_tokens=1536, system=system, tools=TOOLS,
                    messages=history + [{"role": "assistant", "content": response.content},
                    {"role": "user", "content": _build_tool_results(response.content, tool_use.id, tool_result)}])

            elif tool_use.name == "summary_writer":
                import asyncio; loop = asyncio.get_running_loop()
                tool_result = await loop.run_in_executor(None, fetch_summary_writer,
                    tool_use.input["topic"], tool_use.input.get("max_points", 7),
                    tool_use.input.get("lang", "zh-tw"))
                response = client.messages.create(model="claude-sonnet-4-6", max_tokens=1536, system=system, tools=TOOLS,
                    messages=history + [{"role": "assistant", "content": response.content},
                    {"role": "user", "content": _build_tool_results(response.content, tool_use.id, tool_result)}])

            elif tool_use.name == "key_insights":
                import asyncio; loop = asyncio.get_running_loop()
                tool_result = await loop.run_in_executor(None, fetch_key_insights,
                    tool_use.input["topic"], tool_use.input.get("count", 5),
                    tool_use.input.get("lang", "zh-tw"))
                response = client.messages.create(model="claude-sonnet-4-6", max_tokens=1536, system=system, tools=TOOLS,
                    messages=history + [{"role": "assistant", "content": response.content},
                    {"role": "user", "content": _build_tool_results(response.content, tool_use.id, tool_result)}])

            elif tool_use.name == "bias_detector":
                import asyncio; loop = asyncio.get_running_loop()
                tool_result = await loop.run_in_executor(None, fetch_bias_detector,
                    tool_use.input["topic"], tool_use.input.get("lang", "zh-tw"))
                response = client.messages.create(model="claude-sonnet-4-6", max_tokens=1536, system=system, tools=TOOLS,
                    messages=history + [{"role": "assistant", "content": response.content},
                    {"role": "user", "content": _build_tool_results(response.content, tool_use.id, tool_result)}])

            elif tool_use.name == "second_opinion":
                import asyncio; loop = asyncio.get_running_loop()
                tool_result = await loop.run_in_executor(None, fetch_second_opinion,
                    tool_use.input["question"], tool_use.input.get("experts"),
                    tool_use.input.get("lang", "zh-tw"))
                response = client.messages.create(model="claude-sonnet-4-6", max_tokens=2048, system=system, tools=TOOLS,
                    messages=history + [{"role": "assistant", "content": response.content},
                    {"role": "user", "content": _build_tool_results(response.content, tool_use.id, tool_result)}])

            elif tool_use.name == "brainstorm":
                import asyncio; loop = asyncio.get_running_loop()
                tool_result = await loop.run_in_executor(None, fetch_brainstorm,
                    tool_use.input["problem"], tool_use.input.get("count", 8),
                    tool_use.input.get("style", "實用"), tool_use.input.get("lang", "zh-tw"))
                response = client.messages.create(model="claude-sonnet-4-6", max_tokens=2048, system=system, tools=TOOLS,
                    messages=history + [{"role": "assistant", "content": response.content},
                    {"role": "user", "content": _build_tool_results(response.content, tool_use.id, tool_result)}])

            elif tool_use.name == "benchmark_analysis":
                import asyncio; loop = asyncio.get_running_loop()
                tool_result = await loop.run_in_executor(None, fetch_benchmark_analysis,
                    tool_use.input["subject"], tool_use.input.get("industry", ""),
                    tool_use.input.get("lang", "zh-tw"))
                response = client.messages.create(model="claude-sonnet-4-6", max_tokens=2048, system=system, tools=TOOLS,
                    messages=history + [{"role": "assistant", "content": response.content},
                    {"role": "user", "content": _build_tool_results(response.content, tool_use.id, tool_result)}])

            elif tool_use.name == "steel_man":
                import asyncio; loop = asyncio.get_running_loop()
                tool_result = await loop.run_in_executor(None, fetch_steel_man,
                    tool_use.input["opposing_view"], tool_use.input.get("own_position", ""),
                    tool_use.input.get("lang", "zh-tw"))
                response = client.messages.create(model="claude-sonnet-4-6", max_tokens=2048, system=system, tools=TOOLS,
                    messages=history + [{"role": "assistant", "content": response.content},
                    {"role": "user", "content": _build_tool_results(response.content, tool_use.id, tool_result)}])

            elif tool_use.name == "socratic_questioning":
                import asyncio; loop = asyncio.get_running_loop()
                tool_result = await loop.run_in_executor(None, fetch_socratic_questioning,
                    tool_use.input["topic"], tool_use.input.get("depth", 5),
                    tool_use.input.get("lang", "zh-tw"))
                response = client.messages.create(model="claude-sonnet-4-6", max_tokens=2048, system=system, tools=TOOLS,
                    messages=history + [{"role": "assistant", "content": response.content},
                    {"role": "user", "content": _build_tool_results(response.content, tool_use.id, tool_result)}])

            elif tool_use.name == "analogy_maker":
                import asyncio; loop = asyncio.get_running_loop()
                tool_result = await loop.run_in_executor(None, fetch_analogy_maker,
                    tool_use.input["concept"], tool_use.input.get("audience", "一般大眾"),
                    tool_use.input.get("count", 3), tool_use.input.get("lang", "zh-tw"))
                response = client.messages.create(model="claude-sonnet-4-6", max_tokens=1536, system=system, tools=TOOLS,
                    messages=history + [{"role": "assistant", "content": response.content},
                    {"role": "user", "content": _build_tool_results(response.content, tool_use.id, tool_result)}])

            elif tool_use.name == "narrative_builder":
                import asyncio; loop = asyncio.get_running_loop()
                tool_result = await loop.run_in_executor(None, fetch_narrative_builder,
                    tool_use.input["topic"], tool_use.input.get("key_message", ""),
                    tool_use.input.get("audience", ""), tool_use.input.get("lang", "zh-tw"))
                response = client.messages.create(model="claude-sonnet-4-6", max_tokens=2048, system=system, tools=TOOLS,
                    messages=history + [{"role": "assistant", "content": response.content},
                    {"role": "user", "content": _build_tool_results(response.content, tool_use.id, tool_result)}])

            elif tool_use.name == "critique_writer":
                import asyncio; loop = asyncio.get_running_loop()
                tool_result = await loop.run_in_executor(None, fetch_critique_writer,
                    tool_use.input["subject"], tool_use.input.get("type", "觀點"),
                    tool_use.input.get("lang", "zh-tw"))
                response = client.messages.create(model="claude-sonnet-4-6", max_tokens=2048, system=system, tools=TOOLS,
                    messages=history + [{"role": "assistant", "content": response.content},
                    {"role": "user", "content": _build_tool_results(response.content, tool_use.id, tool_result)}])

            elif tool_use.name == "position_statement":
                import asyncio; loop = asyncio.get_running_loop()
                tool_result = await loop.run_in_executor(None, fetch_position_statement,
                    tool_use.input["issue"], tool_use.input["stance"],
                    tool_use.input.get("lang", "zh-tw"))
                response = client.messages.create(model="claude-sonnet-4-6", max_tokens=2048, system=system, tools=TOOLS,
                    messages=history + [{"role": "assistant", "content": response.content},
                    {"role": "user", "content": _build_tool_results(response.content, tool_use.id, tool_result)}])

            elif tool_use.name == "ocr_click":
                import asyncio; loop = asyncio.get_running_loop()
                tool_result = await loop.run_in_executor(None, fetch_ocr_click,
                    tool_use.input["target_text"], tool_use.input.get("monitor", 1),
                    tool_use.input.get("click_type", "click"), tool_use.input.get("region"))
                response = client.messages.create(model="claude-sonnet-4-6", max_tokens=512, system=system, tools=TOOLS,
                    messages=history + [{"role": "assistant", "content": response.content},
                    {"role": "user", "content": _build_tool_results(response.content, tool_use.id, tool_result)}])

            elif tool_use.name == "vision_locate":
                import asyncio; loop = asyncio.get_running_loop()
                tool_result = await loop.run_in_executor(None, fetch_vision_locate,
                    tool_use.input["description"], tool_use.input.get("monitor", 1),
                    tool_use.input.get("action", "click"), tool_use.input.get("region"))
                response = client.messages.create(model="claude-sonnet-4-6", max_tokens=512, system=system, tools=TOOLS,
                    messages=history + [{"role": "assistant", "content": response.content},
                    {"role": "user", "content": _build_tool_results(response.content, tool_use.id, tool_result)}])

            elif tool_use.name == "screen_workflow":
                import asyncio; loop = asyncio.get_running_loop()
                tool_result = await loop.run_in_executor(None, fetch_screen_workflow,
                    tool_use.input["steps"])
                response = client.messages.create(model="claude-sonnet-4-6", max_tokens=1024, system=system, tools=TOOLS,
                    messages=history + [{"role": "assistant", "content": response.content},
                    {"role": "user", "content": _build_tool_results(response.content, tool_use.id, tool_result)}])

            elif tool_use.name == "app_navigator":
                import asyncio, re as _re_nav_guard
                loop = asyncio.get_running_loop()
                # ── 程式碼層級封鎖：群組訊息含 @提及 且無螢幕控制關鍵字 → 直接拒絕 ──
                _screen_kw = ("螢幕", "桌面", "screen", "monitor", "控制", "視窗")
                _has_at = bool(_re_nav_guard.search(r'@\w+', user_text or ""))
                _has_screen = any(k in (user_text or "") for k in _screen_kw)
                if is_group and _has_at and not _has_screen:
                    tool_result = "❌ 禁止使用螢幕控制。這是群組對話，@提及只是表達方式。你必須直接在這個群組聊天室裡用文字回覆，給出股票推薦或回答問題，不要說「已傳給XXX」，直接在群組裡說出你的推薦內容。"
                else:
                    tool_result = await loop.run_in_executor(None, fetch_app_navigator,
                        tool_use.input.get("app", "Telegram"), tool_use.input.get("task", ""),
                        tool_use.input.get("input_text", ""), tool_use.input.get("monitor", 1),
                        tool_use.input.get("contact_name", ""))
                # 把第一輪存入 _nav_tail，讓多工具循環能接續
                _nav_first_asst = {"role": "assistant", "content": response.content}
                _nav_first_result = {"role": "user", "content": _build_tool_results(response.content, tool_use.id, tool_result)}
                response = client.messages.create(model="claude-sonnet-4-6", max_tokens=512, system=system, tools=TOOLS,
                    messages=history + [_nav_first_asst, _nav_first_result])
                # 更新 history 讓後續多工具循環能包含第一輪的對話
                history = history + [_nav_first_asst, _nav_first_result]

            elif tool_use.name == "wait_and_click":
                import asyncio; loop = asyncio.get_running_loop()
                tool_result = await loop.run_in_executor(None, fetch_wait_and_click,
                    tool_use.input["target_text"], tool_use.input.get("timeout", 15),
                    tool_use.input.get("monitor", 1), tool_use.input.get("action_after", "click"))
                response = client.messages.create(model="claude-sonnet-4-6", max_tokens=512, system=system, tools=TOOLS,
                    messages=history + [{"role": "assistant", "content": response.content},
                    {"role": "user", "content": _build_tool_results(response.content, tool_use.id, tool_result)}])

            elif tool_use.name == "drag_drop":
                import asyncio; loop = asyncio.get_running_loop()
                tool_result = await loop.run_in_executor(None, fetch_drag_drop,
                    tool_use.input.get("from_x"), tool_use.input.get("from_y"),
                    tool_use.input.get("to_x"), tool_use.input.get("to_y"),
                    tool_use.input.get("from_text", ""), tool_use.input.get("to_text", ""),
                    tool_use.input.get("monitor", 1), tool_use.input.get("duration", 0.5))
                response = client.messages.create(model="claude-sonnet-4-6", max_tokens=512, system=system, tools=TOOLS,
                    messages=history + [{"role": "assistant", "content": response.content},
                    {"role": "user", "content": _build_tool_results(response.content, tool_use.id, tool_result)}])

            elif tool_use.name == "read_screen":
                import asyncio; loop = asyncio.get_running_loop()
                tool_result = await loop.run_in_executor(None, fetch_read_screen,
                    tool_use.input.get("question", "描述螢幕上有什麼"),
                    tool_use.input.get("monitor", 1))
                response = client.messages.create(model="claude-sonnet-4-6", max_tokens=512, system=system, tools=TOOLS,
                    messages=history + [{"role": "assistant", "content": response.content},
                    {"role": "user", "content": _build_tool_results(response.content, tool_use.id, tool_result)}])

            elif tool_use.name == "scroll_at":
                import asyncio; loop = asyncio.get_running_loop()
                tool_result = await loop.run_in_executor(None, fetch_scroll_at,
                    tool_use.input.get("direction", "down"),
                    tool_use.input.get("amount", 3),
                    tool_use.input.get("x"), tool_use.input.get("y"),
                    tool_use.input.get("monitor", 1),
                    tool_use.input.get("description", ""))
                response = client.messages.create(model="claude-sonnet-4-6", max_tokens=256, system=system, tools=TOOLS,
                    messages=history + [{"role": "assistant", "content": response.content},
                    {"role": "user", "content": _build_tool_results(response.content, tool_use.id, tool_result)}])

            elif tool_use.name == "window_manager":
                import asyncio; loop = asyncio.get_running_loop()
                tool_result = await loop.run_in_executor(None, fetch_window_manager,
                    tool_use.input.get("action", "list"),
                    tool_use.input.get("window_name", ""))
                response = client.messages.create(model="claude-sonnet-4-6", max_tokens=512, system=system, tools=TOOLS,
                    messages=history + [{"role": "assistant", "content": response.content},
                    {"role": "user", "content": _build_tool_results(response.content, tool_use.id, tool_result)}])

            elif tool_use.name == "ptt_search":
                import asyncio
                loop = asyncio.get_running_loop()
                tool_result = await loop.run_in_executor(None, ptt_search,
                    tool_use.input["keyword"],
                    tool_use.input.get("board", "Gossiping"),
                    tool_use.input.get("count", 5))
                response = client.messages.create(
                    model="claude-sonnet-4-6", max_tokens=1536, system=system, tools=TOOLS,
                    messages=history + [
                        {"role": "assistant", "content": response.content},
                        {"role": "user", "content": _build_tool_results(response.content, tool_use.id, tool_result)}
                    ]
                )

            elif tool_use.name == "multi_perspective":
                import asyncio
                loop = asyncio.get_running_loop()
                tool_result = await loop.run_in_executor(None, multi_perspective,
                    tool_use.input["topic"], tool_use.input.get("lang", "zh-tw"))
                response = client.messages.create(
                    model="claude-sonnet-4-6", max_tokens=1536, system=system, tools=TOOLS,
                    messages=history + [
                        {"role": "assistant", "content": response.content},
                        {"role": "user", "content": _build_tool_results(response.content, tool_use.id, tool_result)}
                    ]
                )

            elif tool_use.name == "google_trends":
                import asyncio
                loop = asyncio.get_running_loop()
                tool_result = await loop.run_in_executor(None, fetch_google_trends,
                    tool_use.input["keywords"],
                    tool_use.input.get("timeframe", "today 3-m"),
                    tool_use.input.get("geo", "TW"))
                response = client.messages.create(
                    model="claude-sonnet-4-6", max_tokens=1536, system=system, tools=TOOLS,
                    messages=history + [
                        {"role": "assistant", "content": response.content},
                        {"role": "user", "content": _build_tool_results(response.content, tool_use.id, tool_result)}
                    ]
                )

            elif tool_use.name == "read_webpage":
                import asyncio
                loop = asyncio.get_running_loop()
                tool_result = await loop.run_in_executor(None, read_webpage,
                    tool_use.input["url"], tool_use.input.get("max_chars", 3000))
                response = client.messages.create(
                    model="claude-sonnet-4-6", max_tokens=1536, system=system, tools=TOOLS,
                    messages=history + [
                        {"role": "assistant", "content": response.content},
                        {"role": "user", "content": _build_tool_results(response.content, tool_use.id, tool_result)}
                    ]
                )

            elif tool_use.name == "wikipedia_search":
                import asyncio
                loop = asyncio.get_running_loop()
                tool_result = await loop.run_in_executor(None, wikipedia_search,
                    tool_use.input["query"], tool_use.input.get("lang", "zh"))
                response = client.messages.create(
                    model="claude-sonnet-4-6", max_tokens=1536, system=system, tools=TOOLS,
                    messages=history + [
                        {"role": "assistant", "content": response.content},
                        {"role": "user", "content": _build_tool_results(response.content, tool_use.id, tool_result)}
                    ]
                )

            elif tool_use.name == "news_search":
                import asyncio
                loop = asyncio.get_running_loop()
                tool_result = await loop.run_in_executor(None, search_news,
                    tool_use.input["query"], tool_use.input.get("lang", "zh-TW"),
                    tool_use.input.get("count", 6))
                response = client.messages.create(
                    model="claude-sonnet-4-6", max_tokens=1536, system=system, tools=TOOLS,
                    messages=history + [
                        {"role": "assistant", "content": response.content},
                        {"role": "user", "content": _build_tool_results(response.content, tool_use.id, tool_result)}
                    ]
                )

            elif tool_use.name == "youtube_summary":
                import asyncio
                loop = asyncio.get_running_loop()
                tool_result = await loop.run_in_executor(None, youtube_summary, tool_use.input["url"])
                response = client.messages.create(
                    model="claude-sonnet-4-6", max_tokens=1536, system=system, tools=TOOLS,
                    messages=history + [
                        {"role": "assistant", "content": response.content},
                        {"role": "user", "content": _build_tool_results(response.content, tool_use.id, tool_result)}
                    ]
                )

            elif tool_use.name == "analyze_pdf":
                import asyncio
                loop = asyncio.get_running_loop()
                tool_result = await loop.run_in_executor(None, analyze_pdf,
                    tool_use.input["path"], tool_use.input.get("max_chars", 4000))
                response = client.messages.create(
                    model="claude-sonnet-4-6", max_tokens=1536, system=system, tools=TOOLS,
                    messages=history + [
                        {"role": "assistant", "content": response.content},
                        {"role": "user", "content": _build_tool_results(response.content, tool_use.id, tool_result)}
                    ]
                )

            elif tool_use.name == "ddg_search":
                import asyncio
                loop = asyncio.get_running_loop()
                tool_result = await loop.run_in_executor(None, execute_ddg_search,
                    tool_use.input["query"],
                    tool_use.input.get("region", "zh-tw"),
                    tool_use.input.get("max_results", 5))
                response = client.messages.create(
                    model="claude-sonnet-4-6", max_tokens=1024, system=system, tools=TOOLS,
                    messages=history + [
                        {"role": "assistant", "content": response.content},
                        {"role": "user", "content": _build_tool_results(response.content, tool_use.id, tool_result)}
                    ]
                )

            elif tool_use.name == "get_fundamentals":
                import asyncio
                loop = asyncio.get_running_loop()
                tool_result = await loop.run_in_executor(None, fetch_fundamentals, tool_use.input["symbol"])
                response = client.messages.create(
                    model="claude-sonnet-4-6", max_tokens=1024, system=system, tools=TOOLS,
                    messages=history + [
                        {"role": "assistant", "content": response.content},
                        {"role": "user", "content": _build_tool_results(response.content, tool_use.id, tool_result)}
                    ]
                )

            elif tool_use.name == "get_market_sentiment":
                import asyncio
                loop = asyncio.get_running_loop()
                tool_result = await loop.run_in_executor(None, fetch_market_sentiment)
                response = client.messages.create(
                    model="claude-sonnet-4-6", max_tokens=1024, system=system, tools=TOOLS,
                    messages=history + [
                        {"role": "assistant", "content": response.content},
                        {"role": "user", "content": _build_tool_results(response.content, tool_use.id, tool_result)}
                    ]
                )

            elif tool_use.name == "get_etf":
                import asyncio
                loop = asyncio.get_running_loop()
                tool_result = await loop.run_in_executor(None, fetch_etf, tool_use.input["symbol"])
                response = client.messages.create(
                    model="claude-sonnet-4-6", max_tokens=1024, system=system, tools=TOOLS,
                    messages=history + [
                        {"role": "assistant", "content": response.content},
                        {"role": "user", "content": _build_tool_results(response.content, tool_use.id, tool_result)}
                    ]
                )

            elif tool_use.name == "get_earnings":
                import asyncio
                loop = asyncio.get_running_loop()
                tool_result = await loop.run_in_executor(None, fetch_earnings, tool_use.input["symbol"])
                response = client.messages.create(
                    model="claude-sonnet-4-6", max_tokens=1024, system=system, tools=TOOLS,
                    messages=history + [
                        {"role": "assistant", "content": response.content},
                        {"role": "user", "content": _build_tool_results(response.content, tool_use.id, tool_result)}
                    ]
                )

            elif tool_use.name == "get_candlestick_chart":
                import asyncio
                loop = asyncio.get_running_loop()
                img_bytes, pattern_str = await loop.run_in_executor(None, generate_candlestick,
                    tool_use.input["symbol"], tool_use.input.get("period", "3mo"))
                if img_bytes:
                    import io
                    await update.message.reply_photo(photo=io.BytesIO(img_bytes), caption=f"型態分析：\n{pattern_str}")
                tool_result = pattern_str
                response = client.messages.create(
                    model="claude-sonnet-4-6", max_tokens=1024, system=system, tools=TOOLS,
                    messages=history + [
                        {"role": "assistant", "content": response.content},
                        {"role": "user", "content": _build_tool_results(response.content, tool_use.id, tool_result)}
                    ]
                )

            elif tool_use.name == "compare_stocks":
                import asyncio
                loop = asyncio.get_running_loop()
                tool_result = await loop.run_in_executor(None, compare_stocks,
                    tool_use.input["symbols"], tool_use.input.get("metrics"))
                response = client.messages.create(
                    model="claude-sonnet-4-6", max_tokens=1024, system=system, tools=TOOLS,
                    messages=history + [
                        {"role": "assistant", "content": response.content},
                        {"role": "user", "content": _build_tool_results(response.content, tool_use.id, tool_result)}
                    ]
                )

            elif tool_use.name == "get_stock_advanced":
                import asyncio
                loop = asyncio.get_running_loop()
                tool_result = await loop.run_in_executor(None, fetch_stock_advanced,
                    tool_use.input["symbol"], tool_use.input.get("indicators"))
                response = client.messages.create(
                    model="claude-sonnet-4-6", max_tokens=1024, system=system, tools=TOOLS,
                    messages=history + [
                        {"role": "assistant", "content": response.content},
                        {"role": "user", "content": _build_tool_results(response.content, tool_use.id, tool_result)}
                    ]
                )

            elif tool_use.name == "get_macro":
                import asyncio
                loop = asyncio.get_running_loop()
                tool_result = await loop.run_in_executor(None, fetch_macro, tool_use.input["indicator"])
                response = client.messages.create(
                    model="claude-sonnet-4-6", max_tokens=1024, system=system, tools=TOOLS,
                    messages=history + [
                        {"role": "assistant", "content": response.content},
                        {"role": "user", "content": _build_tool_results(response.content, tool_use.id, tool_result)}
                    ]
                )

            elif tool_use.name == "portfolio":
                import asyncio
                loop = asyncio.get_running_loop()
                tool_result = await loop.run_in_executor(None, execute_portfolio,
                    tool_use.input["action"],
                    tool_use.input.get("chat_id", chat_id),
                    tool_use.input.get("symbol", ""),
                    tool_use.input.get("shares", 0),
                    tool_use.input.get("cost", 0)
                )
                response = client.messages.create(
                    model="claude-sonnet-4-6", max_tokens=1024, system=system, tools=TOOLS,
                    messages=history + [
                        {"role": "assistant", "content": response.content},
                        {"role": "user", "content": _build_tool_results(response.content, tool_use.id, tool_result)}
                    ]
                )

            elif tool_use.name == "desktop_control":
                import asyncio as _asyncio2
                _loop2 = _asyncio2.get_running_loop()
                _dc_messages = list(history)
                _dc_cur_resp = response
                _dc_cur_tu = tool_use
                _dc_last_msg = ""
                _dc_last_screenshot = None

                # 統一迴圈：desktop_control 和 simple_tools 都在同一層處理
                for _dc_step in range(15):
                    _dc_tool_result = ""

                    if _dc_cur_tu.name == "desktop_control":
                        _dc_inp = _dc_cur_tu.input
                        # ── 攔截：screenshot 但實際是 Telegram/LINE 導航請求 ──
                        import re as _re_nav
                        _is_tg_nav = (
                            _dc_inp.get("action") == "screenshot" and
                            _re_nav.search(r'telegram|line', user_text, _re_nav.IGNORECASE) and
                            _re_nav.search(r'找.{1,20}(說|傳|回覆|發|跟他說)', user_text)
                        )
                        if _is_tg_nav:
                            _mon_m = _re_nav.search(r'螢幕(\d)', user_text)
                            _nav_mon = int(_mon_m.group(1)) if _mon_m else 1
                            _nav_app = "LINE" if _re_nav.search(r'line', user_text, _re_nav.IGNORECASE) else "Telegram"
                            _nm = _re_nav.search(r'找(.+?)(?:跟他說|並說|跟.{0,5}說|說|給|傳|發)(.+)', user_text)
                            if _nm:
                                _nav_task = _nm.group(1).strip()
                                _nav_input = _nm.group(2).strip()
                            else:
                                _nav_task = user_text
                                _nav_input = ""
                            _nav_r = await _loop2.run_in_executor(
                                None, fetch_app_navigator, _nav_app, _nav_task, _nav_input, _nav_mon, _nav_task
                            )
                            _dc_result = {"ok": True, "message": _nav_r, "screenshot": None}
                        else:
                            _dc_result = await _loop2.run_in_executor(
                                None,
                                lambda _i=_dc_inp: execute_desktop_control(
                                    action=_i["action"],
                                    x=_i.get("x"),
                                    y=_i.get("y"),
                                    text=_i.get("text"),
                                    app=_i.get("app"),
                                    direction=_i.get("direction", "down"),
                                    amount=_i.get("amount", 3),
                                    monitor=_i.get("monitor")
                                )
                            )
                        if _dc_result.get("screenshot"):
                            _dc_last_screenshot = _dc_result["screenshot"]
                        _dc_last_msg = _dc_result["message"]
                        _dc_tool_result = _dc_last_msg
                        if _dc_inp.get("action") == "open_app":
                            await _asyncio2.sleep(1.0)
                    else:
                        # simple_tools
                        tool_use = _dc_cur_tu
                        _dc_tool_result = await _loop2.run_in_executor(None, simple_tools[tool_use.name])
                        _dc_last_msg = _dc_tool_result


                    _dc_messages = _dc_messages + [
                        {"role": "assistant", "content": _dc_cur_resp.content},
                        {"role": "user", "content": _build_tool_results(_dc_cur_resp.content, _dc_cur_tu.id, _dc_tool_result)}
                    ]
                    _dc_next_resp = client.messages.create(
                        model="claude-sonnet-4-6",
                        max_tokens=512,
                        system=system,
                        tools=TOOLS,
                        messages=_dc_messages
                    )

                    if _dc_next_resp.stop_reason == "tool_use":
                        _dc_next_tu = next((b for b in _dc_next_resp.content if b.type == "tool_use"), None)
                        if _dc_next_tu and (_dc_next_tu.name == "desktop_control" or _dc_next_tu.name in simple_tools):
                            _dc_cur_tu = _dc_next_tu
                            _dc_cur_resp = _dc_next_resp
                            continue

                    _dc_cur_resp = _dc_next_resp
                    break

                # 如果有截圖，先傳圖（失敗不中斷）
                if _dc_last_screenshot:
                    try:
                        await update.message.reply_photo(photo=_dc_last_screenshot, caption="📸 桌面截圖")
                    except Exception:
                        pass
                text_blocks = [b.text for b in _dc_cur_resp.content if hasattr(b, "text")]
                reply = text_blocks[0] if text_blocks else _dc_last_msg
                save_message(chat_id, "assistant", reply)
                await _fr(reply)
                return

            elif tool_use.name == "manage_schedule":
                inp = tool_use.input
                import asyncio
                loop = asyncio.get_running_loop()
                tool_result = await loop.run_in_executor(
                    None,
                    lambda: execute_manage_schedule(
                        action=inp["action"],
                        name=inp.get("name", ""),
                        time=inp.get("time", ""),
                        script=inp.get("script", "")
                    )
                )
                response = client.messages.create(
                    model="claude-sonnet-4-6",
                    max_tokens=512,
                    system=system,
                    tools=TOOLS,
                    messages=history + [
                        {"role": "assistant", "content": response.content},
                        {"role": "user", "content": _build_tool_results(response.content, tool_use.id, tool_result)}
                    ]
                )
                text_blocks = [b.text for b in response.content if hasattr(b, "text")]
                reply = text_blocks[0] if text_blocks else tool_result
                save_message(chat_id, "assistant", reply)
                await _fr(reply)
                return

            elif tool_use.name == "generate_image":
                prompt = tool_use.input["prompt"]
                width = tool_use.input.get("width", 512)
                height = tool_use.input.get("height", 512)
                overlay_text = tool_use.input.get("overlay_text", "")
                await context.bot.send_chat_action(chat_id=chat_id, action="upload_photo")
                import asyncio
                loop = asyncio.get_running_loop()
                image_bytes = await loop.run_in_executor(None, fetch_image, prompt, width, height)
                if image_bytes:
                    if overlay_text:
                        image_bytes = add_text_overlay(image_bytes, overlay_text)
                if image_bytes:
                    await update.message.reply_photo(photo=image_bytes, caption=f"🎨 {prompt}")
                else:
                    await update.message.reply_text("圖片生成失敗，請稍後再試。")
                save_message(chat_id, "assistant", f"已為用戶生成圖片：{prompt}")
                return

        # ── 多輪工具 dispatch：處理連續/平行 tool calls ──
        # 用 _tail 追蹤本輪所有 tool 交換，確保 messages 格式正確
        import asyncio as _aio
        _eloop = _aio.get_running_loop()
        _tail = []        # 本輪累積的 [assistant, user] pairs
        _extra_rounds = 0

        def _tool_dispatch(tu):
            """根據 tool 名稱呼叫對應 fetch 函式"""
            n, i = tu.name, tu.input
            _map = {
                "get_stock":          lambda: fetch_stock(i.get("symbol",""), i.get("period","1mo")),
                "get_fundamentals":   lambda: fetch_fundamentals(i.get("symbol","")),
                "get_market_sentiment": lambda: fetch_market_sentiment(),
                "get_sector":         lambda: fetch_sector(i.get("market","us")),
                "get_global_market":  lambda: fetch_global_market(),
                "get_weather":        lambda: fetch_weather(i.get("city","台北")),
                "get_crypto":         lambda: fetch_crypto(i.get("coin","btc"), i.get("vs_currency","usd")),
                "get_ashare":         lambda: fetch_ashare(i.get("code",""), i.get("period","1mo")),
                "compare_stocks":     lambda: compare_stocks(i.get("symbols",[]), i.get("metrics",["all"])),
                "stock_screener":     lambda: fetch_stock_screener(i.get("criteria",""), i.get("market","tw")),
                "get_analyst_ratings": lambda: fetch_analyst_ratings(i.get("symbol","")),
                "get_risk_metrics":   lambda: fetch_risk_metrics(i.get("symbol",""), i.get("period","1y")),
                "get_earnings":       lambda: fetch_earnings(i.get("symbol","")),
                "get_finance_news":   lambda: fetch_finance_news(i.get("source","all"), i.get("count",5)),
                "get_etf":            lambda: fetch_etf(i.get("symbol","")),
                "get_earnings_calendar": lambda: fetch_earnings_calendar(i.get("days",7)),
                "get_dividend_calendar": lambda: fetch_dividend_calendar(i.get("symbol","")),
                "ptt_search":         lambda: ptt_search(i.get("keyword",""), i.get("board","Gossiping"), i.get("count",5)),
                "china_search":       lambda: fetch_china_search(i.get("query",""), i.get("category","其他"), i.get("count",6)),
                "read_screen":        lambda: fetch_read_screen(i.get("question","描述螢幕上有什麼"), i.get("monitor",1)),
                "ocr_click":          lambda: fetch_ocr_click(i.get("target_text",""), i.get("monitor",1), i.get("click_type","click"), i.get("region")),
                "vision_locate":      lambda: fetch_vision_locate(i.get("description",""), i.get("monitor",1), i.get("action","click"), i.get("region")),
                "scroll_at":          lambda: fetch_scroll_at(i.get("direction","down"), i.get("amount",3), i.get("x"), i.get("y"), i.get("monitor",1), i.get("description","")),
                "window_manager":     lambda: fetch_window_manager(i.get("action","list"), i.get("window_name","")),
                "app_navigator":      lambda: fetch_app_navigator(i.get("app","Telegram"), i.get("task",""), i.get("input_text",""), i.get("monitor",1), i.get("contact_name","")),
                # ── 金融工具 ──
                "get_forex":              lambda: fetch_forex(i.get("pair", "")),
                "get_cn_news":            lambda: fetch_cn_news(i.get("source", "all"), i.get("count", 5)),
                "get_commodity":          lambda: fetch_commodity(i.get("items", ["all"])),
                "get_bond_yield":         lambda: fetch_bond_yield(),
                "get_margin_trading":     lambda: fetch_margin_trading(i.get("symbol", ""), i.get("date", "")),
                "get_options":            lambda: fetch_options(i.get("symbol", ""), i.get("expiry", "")),
                "get_futures":            lambda: fetch_futures(i.get("items", ["all"])),
                "get_ipo":                lambda: fetch_ipo(i.get("count", 10)),
                "backtest":               lambda: fetch_backtest(i.get("symbol", ""), i.get("strategy", "ma_cross"), i.get("period", "2y")),
                "get_economic_calendar":  lambda: fetch_economic_calendar(i.get("count", 10)),
                "get_short_interest":     lambda: fetch_short_interest(i.get("symbol", "")),
                "get_correlation":        lambda: fetch_correlation(i.get("symbols", []), i.get("period", "1y")),
                "get_money_flow":         lambda: fetch_money_flow(i.get("symbol", "")),
                "get_concept_stocks":     lambda: fetch_concept_stocks(i.get("theme", "")),
                "get_crypto_depth":       lambda: fetch_crypto_depth(i.get("coin", "bitcoin")),
                "drip_calculator":        lambda: fetch_drip_calculator(i.get("symbol", ""), i.get("shares", 0), i.get("years", 10), i.get("monthly_invest", 0)),
                "get_forex_chart":        lambda: fetch_forex_chart(i.get("pair", ""), i.get("period", "3mo")),
                "get_warrant":            lambda: fetch_warrant(i.get("underlying", "")),
                "get_portfolio_risk":     lambda: fetch_portfolio_risk(i.get("holdings", []), i.get("period", "1y")),
                "get_institutional":      lambda: fetch_institutional(i.get("symbol", ""), i.get("days", 5)),
                "retirement_calculator":  lambda: fetch_retirement_calculator(i.get("current_age", 30), i.get("current_savings", 0), i.get("monthly_save", 0), i.get("retire_age", 65), i.get("annual_return", 6.0), i.get("monthly_expense", 50000)),
                "loan_calculator":        lambda: fetch_loan_calculator(i.get("principal", 0), i.get("annual_rate", 0), i.get("years", 0), i.get("loan_type", "等額本息")),
                "compound_calculator":    lambda: fetch_compound_calculator(i.get("principal", 0), i.get("annual_rate", 0), i.get("years", 0), i.get("monthly_add", 0), i.get("compound_freq", 12)),
                "asset_allocation":       lambda: fetch_asset_allocation(i.get("age", 30), i.get("risk_level", ""), i.get("goal", "退休"), i.get("investment_horizon")),
                "tw_tax_calculator":      lambda: fetch_tw_tax_calculator(i.get("dividend_income", 0), i.get("other_income", 0), i.get("tax_bracket"), i.get("sell_amount", 0)),
                "currency_converter":     lambda: fetch_currency_converter(i.get("amount", 0), i.get("from_currency", ""), i.get("to_currency", "")),
                "get_fund":               lambda: fetch_fund(i.get("symbol", "")),
                "get_reits":              lambda: fetch_reits(i.get("symbol", "")),
                "inflation_adjusted":     lambda: fetch_inflation_adjusted(i.get("nominal_return", 0), i.get("years", 0), i.get("amount", 0), i.get("inflation_rate", 2.0)),
                "defi_calculator":        lambda: fetch_defi_calculator(i.get("principal_usd", 0), i.get("apy", 0), i.get("days", 0), i.get("compound", True), i.get("protocol", "")),
                "gold_calculator":        lambda: fetch_gold_calculator(i.get("weight", 0), i.get("unit", "公克"), i.get("currency", "TWD")),
                "forex_deposit":          lambda: fetch_forex_deposit(i.get("amount_twd", 0), i.get("currency", ""), i.get("annual_rate", 0), i.get("months", 0), i.get("buy_rate"), i.get("sell_rate")),
                "financial_health":       lambda: fetch_financial_health(i.get("monthly_income", 0), i.get("monthly_expense", 0), i.get("total_assets", 0), i.get("total_debt", 0), i.get("emergency_fund_months", 0), i.get("has_insurance", False), i.get("investment_ratio", 0)),
                "get_candlestick_chart":  lambda: generate_candlestick(i.get("symbol", ""), i.get("period", "3mo"))[1],
                "get_stock_advanced":     lambda: fetch_stock_advanced(i.get("symbol", ""), i.get("indicators")),
                "get_macro":              lambda: fetch_macro(i.get("indicator", "")),
                "portfolio":              lambda: execute_portfolio(i.get("action", ""), i.get("chat_id") or chat_id, i.get("symbol", ""), i.get("shares", 0), i.get("cost", 0)),
                # ── 研究工具 ──
                "deep_research":          lambda: fetch_deep_research(i.get("topic", ""), i.get("lang", "zh-tw"), i.get("depth", 5)),
                "fact_check":             lambda: fetch_fact_check(i.get("claim", ""), i.get("lang", "zh-tw")),
                "timeline_events":        lambda: fetch_timeline_events(i.get("topic", ""), i.get("lang", "zh-tw")),
                "sentiment_scan":         lambda: fetch_sentiment_scan(i.get("topic", ""), i.get("lang", "zh-tw")),
                "compare_analysis":       lambda: fetch_compare_analysis(i.get("items", []), i.get("dimensions"), i.get("context", "")),
                "pros_cons_analysis":     lambda: fetch_pros_cons_analysis(i.get("subject", ""), i.get("context", ""), i.get("lang", "zh-tw")),
                "research_report":        lambda: fetch_research_report(i.get("topic", ""), i.get("purpose", "一般研究"), i.get("lang", "zh-tw")),
                "opinion_writer":         lambda: fetch_opinion_writer(i.get("topic", ""), i.get("stance", "中立"), i.get("style", "正式")),
                "trend_forecast":         lambda: fetch_trend_forecast(i.get("topic", ""), i.get("timeframe", "全部"), i.get("lang", "zh-tw")),
                "debate_simulator":       lambda: fetch_debate_simulator(i.get("motion", ""), i.get("lang", "zh-tw")),
                "academic_search":        lambda: fetch_academic_search(i.get("query", ""), i.get("field", ""), i.get("lang", "en")),
                "health_research":        lambda: fetch_health_research(i.get("topic", ""), i.get("lang", "zh-tw")),
                "law_research":           lambda: fetch_law_research(i.get("topic", ""), i.get("jurisdiction", "台灣"), i.get("lang", "zh-tw")),
                "person_research":        lambda: fetch_person_research(i.get("name", ""), i.get("context", ""), i.get("lang", "zh-tw")),
                "company_research":       lambda: fetch_company_research(i.get("company", ""), i.get("lang", "zh-tw")),
                "product_review":         lambda: fetch_product_review(i.get("product", ""), i.get("category", ""), i.get("lang", "zh-tw")),
                "travel_research":        lambda: fetch_travel_research(i.get("destination", ""), i.get("days"), i.get("style", ""), i.get("lang", "zh-tw")),
                "job_market":             lambda: fetch_job_market(i.get("job_title", ""), i.get("location", "台灣"), i.get("lang", "zh-tw")),
                "impact_analysis":        lambda: fetch_impact_analysis(i.get("event", ""), i.get("scope"), i.get("lang", "zh-tw")),
                "scenario_planning":      lambda: fetch_scenario_planning(i.get("topic", ""), i.get("horizon", ""), i.get("lang", "zh-tw")),
                "decision_helper":        lambda: fetch_decision_helper(i.get("question", ""), i.get("options"), i.get("criteria")),
                "devil_advocate":         lambda: fetch_devil_advocate(i.get("position", ""), i.get("lang", "zh-tw")),
                "summary_writer":         lambda: fetch_summary_writer(i.get("topic", ""), i.get("max_points", 7), i.get("lang", "zh-tw")),
                "key_insights":           lambda: fetch_key_insights(i.get("topic", ""), i.get("count", 5), i.get("lang", "zh-tw")),
                "bias_detector":          lambda: fetch_bias_detector(i.get("topic", ""), i.get("lang", "zh-tw")),
                "second_opinion":         lambda: fetch_second_opinion(i.get("question", ""), i.get("experts"), i.get("lang", "zh-tw")),
                "brainstorm":             lambda: fetch_brainstorm(i.get("problem", ""), i.get("count", 8), i.get("style", "實用"), i.get("lang", "zh-tw")),
                "benchmark_analysis":     lambda: fetch_benchmark_analysis(i.get("subject", ""), i.get("industry", ""), i.get("lang", "zh-tw")),
                "steel_man":              lambda: fetch_steel_man(i.get("opposing_view", ""), i.get("own_position", ""), i.get("lang", "zh-tw")),
                "socratic_questioning":   lambda: fetch_socratic_questioning(i.get("topic", ""), i.get("depth", 5), i.get("lang", "zh-tw")),
                "analogy_maker":          lambda: fetch_analogy_maker(i.get("concept", ""), i.get("audience", "一般大眾"), i.get("count", 3), i.get("lang", "zh-tw")),
                "narrative_builder":      lambda: fetch_narrative_builder(i.get("topic", ""), i.get("key_message", ""), i.get("audience", ""), i.get("lang", "zh-tw")),
                "critique_writer":        lambda: fetch_critique_writer(i.get("subject", ""), i.get("type", "觀點"), i.get("lang", "zh-tw")),
                "position_statement":     lambda: fetch_position_statement(i.get("issue", ""), i.get("stance", ""), i.get("lang", "zh-tw")),
                # ── 搜尋/網頁工具 ──
                "ddg_search":             lambda: execute_ddg_search(i.get("query", ""), i.get("region", "zh-tw"), i.get("max_results", 5)),
                "multi_perspective":      lambda: multi_perspective(i.get("topic", ""), i.get("lang", "zh-tw")),
                "google_trends":          lambda: fetch_google_trends(i.get("keywords", []), i.get("timeframe", "today 3-m"), i.get("geo", "TW")),
                "read_webpage":           lambda: read_webpage(i.get("url", ""), i.get("max_chars", 3000)),
                "wikipedia_search":       lambda: wikipedia_search(i.get("query", ""), i.get("lang", "zh")),
                "news_search":            lambda: search_news(i.get("query", ""), i.get("lang", "zh-TW"), i.get("count", 6)),
                "youtube_summary":        lambda: youtube_summary(i.get("url", "")),
                "analyze_pdf":            lambda: analyze_pdf(i.get("path", ""), i.get("max_chars", 4000)),
                # ── 螢幕/桌面工具 ──
                "screen_workflow":        lambda: fetch_screen_workflow(i.get("steps", [])),
                "wait_and_click":         lambda: fetch_wait_and_click(i.get("target_text", ""), i.get("timeout", 15), i.get("monitor", 1), i.get("action_after", "click")),
                "drag_drop":              lambda: fetch_drag_drop(i.get("from_x"), i.get("from_y"), i.get("to_x"), i.get("to_y"), i.get("from_text", ""), i.get("to_text", ""), i.get("monitor", 1), i.get("duration", 0.5)),
            }
            fn = _map.get(n)
            return fn() if fn else f"工具 {n} 已執行（無對應 handler）"

        while True:
            _has_tool_use = any(hasattr(b, "type") and b.type == "tool_use" for b in response.content)
            text_blocks = [b.text for b in response.content if hasattr(b, "text")]
            # 只在「有 text 且沒有 tool_use」時才停，有 tool_use 就繼續執行
            if text_blocks and not _has_tool_use:
                break
            if not _has_tool_use or _extra_rounds >= 15:
                break
            # 取出本次回應的「所有」tool_use block（含平行呼叫）
            _cur_tus = [b for b in response.content if hasattr(b, "type") and b.type == "tool_use"]
            if not _cur_tus:
                break
            _extra_rounds += 1
            # 執行全部 tool，收集對應 tool_result
            _results_block = []
            for _tu in _cur_tus:
                _res = await _eloop.run_in_executor(None, _tool_dispatch, _tu)
                _results_block.append({"type": "tool_result", "tool_use_id": _tu.id, "content": str(_res)})
            # 累積本輪 assistant + user
            _tail.append({"role": "assistant", "content": response.content})
            _tail.append({"role": "user",      "content": _results_block})
            response = client.messages.create(
                model="claude-sonnet-4-6", max_tokens=1024, system=system, tools=TOOLS,
                messages=history + _tail
            )

        text_blocks = [b.text for b in response.content if hasattr(b, "text")]
        if not text_blocks:
            await update.message.reply_text("哎，我卡住了，請再說一次。")
            return
        reply = text_blocks[0]

        # ── 最終攔截：回覆裡說「幫你點」或用戶原始指令有「播放」→ 強制 vision_locate ──
        import re as _re_final
        _user_said_play = bool(_re_final.search(r'撥放|播放', user_text or ""))
        if _re_final.search(r'幫你點第一|幫你點|現在.*點第一|點第一首', reply) or _user_said_play:
            _vl_final = await loop.run_in_executor(
                None, fetch_vision_locate, "YouTube搜尋結果中第一個非廣告的影片縮圖", 2, "click"
            )
            if "找不到" in str(_vl_final):
                _vl_final = await loop.run_in_executor(
                    None, fetch_vision_locate, "YouTube搜尋結果中第一個非廣告的影片縮圖", 1, "click"
                )
            _s = sender_name if not is_owner else "于晏哥"
            if "找到" in str(_vl_final) and "找不到" not in str(_vl_final):
                reply = f"點好了{_s}！！影片開始播了！！🎵🐮🐴"
            else:
                reply = f"找不到影片耶{_s}，螢幕上可能沒有YouTube搜尋結果頁面😅🐮🐴"

        save_message(chat_id, "assistant", reply)
        log_message("<<", "小牛馬", chat_id, reply)
        # 若 streaming 已送出部分訊息，直接 edit 最終完整內容
        if sent_msg is not None:
            try:
                final_text = reply if is_owner else _fix_group_reply(reply, sender_name)
                await sent_msg.edit_text(final_text)
            except Exception as _edit_err:
                # "Message is not modified" 代表內容一樣，不需要再送一則，其餘錯誤才補發
                if "not modified" not in str(_edit_err).lower():
                    await _fr(reply)
        else:
            await _fr(reply)

    except Exception as e:
        err_str = str(e)
        logging.error(f"handle_message error: {e}", exc_info=True)
        # Timeout 錯誤：通知用戶但不 crash
        if "timed out" in err_str.lower() or "timeout" in err_str.lower():
            try:
                await update.message.reply_text("⏱️ 操作太久逾時了，但動作可能已經執行完了，你看一下螢幕結果。")
            except Exception:
                pass
        # tool_use 沒有 tool_result 的 400 錯誤：清歷史重試提示
        elif "tool_use" in err_str and "tool_result" in err_str:
            try:
                await update.message.reply_text("🔄 對話紀錄出了點問題，用 /clear 清一下再試。")
            except Exception:
                pass
        else:
            try:
                await update.message.reply_text("發生錯誤，請稍後再試。")
            except Exception:
                pass


async def chatid(update: Update, context: ContextTypes.DEFAULT_TYPE):
    await update.message.reply_text(f"這個聊天室的 ID 是：`{update.effective_chat.id}`", parse_mode="Markdown")

async def myid(update: Update, context: ContextTypes.DEFAULT_TYPE):
    await update.message.reply_text(f"你的 Telegram ID 是：`{update.effective_user.id}`", parse_mode="Markdown")


async def clear(update: Update, context: ContextTypes.DEFAULT_TYPE):
    clear_history(update.effective_chat.id)
    await update.message.reply_text("對話紀錄已清除。")


async def memories(update: Update, context: ContextTypes.DEFAULT_TYPE):
    mems = load_long_term_memory(update.effective_chat.id)
    if not mems:
        await update.message.reply_text("目前沒有長期記憶。")
        return
    text = "📝 長期記憶列表：\n\n" + "\n".join(f"[{m['id']}] {m['content']}" for m in mems)
    await update.message.reply_text(text)


async def goodmorning(context: ContextTypes.DEFAULT_TYPE):
    logging.info("【排程觸發】goodmorning 開始執行")
    with open("C:/Users/blue_/claude-telegram-bot/schedule.log", "a", encoding="utf-8") as f:
        f.write(f"[{datetime.datetime.now().strftime('%Y-%m-%d %H:%M:%S')}] goodmorning 觸發\n")
    group_id = int(os.getenv("GROUP_CHAT_ID"))
    import asyncio
    loop = asyncio.get_running_loop()

    # 抓今日世界重要新聞
    def fetch_world_news():
        try:
            import feedparser
            feeds = [
                "https://news.google.com/rss/topics/CAAqJggKIiBDQkFTRWdvSUwyMHZNRGx1YlY4U0FucG9HZVFCQVAB?hl=zh-TW&gl=TW&ceid=TW:zh-Hant",
                "https://feeds.bbci.co.uk/news/world/rss.xml",
            ]
            headlines = []
            for url in feeds:
                try:
                    feed = feedparser.parse(url)
                    for entry in feed.entries[:5]:
                        headlines.append(entry.get("title", ""))
                except Exception:
                    pass
            return "\n".join(headlines[:10]) if headlines else "無法取得新聞"
        except ImportError:
            # feedparser 不可用，用 requests + BeautifulSoup
            try:
                from bs4 import BeautifulSoup
                res = requests.get("https://news.google.com/topics/CAAqJggKIiBDQkFTRWdvSUwyMHZNRGx1YlY4U0FucG9HZVFCQVAB?hl=zh-TW&gl=TW&ceid=TW:zh-Hant", timeout=10, headers={"User-Agent": "Mozilla/5.0"})
                soup = BeautifulSoup(res.text, "html.parser")
                headlines = [a.get_text(strip=True) for a in soup.select("article h3 a, article h4 a")[:10]]
                return "\n".join(headlines) if headlines else "無法取得新聞"
            except Exception as e:
                return f"新聞抓取失敗：{e}"

    news_raw = await loop.run_in_executor(None, fetch_world_news)

    def generate_goodmorning():
        response = client.messages.create(
            model="claude-sonnet-4-6",
            max_tokens=400,
            system=(
                "你是小牛馬，語氣像周杰倫那樣低調自然，帶一點溫暖。現在是早上11點，用繁體中文生成早安訊息給員工看，包含兩部分：\n"
                "1. 一段溫馨勵志的早安問候（30-60字），讓人看了有動力開始一天，但不要太正式或說教，用自然的台灣口語說出來。\n"
                "2. 今日世界重要事件摘要：從提供的新聞標題中挑出 3～5 則最重要的，用一句話簡述每則，輕鬆帶過。\n"
                "整體風格溫馨但不做作，像一個可靠的夥伴跟大家說早安。"
            ),
            messages=[{"role": "user", "content": f"生成今天的早安訊息。\n\n今日新聞標題：\n{news_raw}"}]
        )
        return response.content[0].text
    try:
        msg = await loop.run_in_executor(None, generate_goodmorning)
        await context.bot.send_message(chat_id=group_id, text=msg)
    except Exception as e:
        logging.error(f"goodmorning 發送失敗：{e}", exc_info=True)

async def goodnight(context: ContextTypes.DEFAULT_TYPE):
    logging.info("【排程觸發】goodnight 開始執行")
    with open("C:/Users/blue_/claude-telegram-bot/schedule.log", "a", encoding="utf-8") as f:
        f.write(f"[{datetime.datetime.now().strftime('%Y-%m-%d %H:%M:%S')}] goodnight 觸發\n")
    group_id = int(os.getenv("GROUP_CHAT_ID"))
    import asyncio
    loop = asyncio.get_running_loop()
    def generate_goodnight():
        response = client.messages.create(
            model="claude-sonnet-4-6",
            max_tokens=150,
            system="你是小牛馬，語氣像周杰倫那樣低調自然，帶一點溫暖。現在是晚上10:30，用繁體中文生成一段溫馨勵志的晚安問候語給員工看，每次內容都不同，讓人看了覺得今天辛苦了、明天繼續加油。用自然的台灣口語，不要說教不要太正式，像一個可靠的夥伴跟大家說晚安。結尾加上表情符號。約30-60字。",
            messages=[{"role": "user", "content": "生成今晚的晚安訊息"}]
        )
        return response.content[0].text
    try:
        msg = await loop.run_in_executor(None, generate_goodnight)
        await context.bot.send_message(chat_id=group_id, text=msg)
    except Exception as e:
        logging.error(f"goodnight 發送失敗：{e}", exc_info=True)

_report_today: str = ""  # 防止同一天重複執行

async def daily_skill_report(context: ContextTypes.DEFAULT_TYPE):
    """每天晚上 9 點：回報今天總共安裝了多少技能"""
    logging.info("【排程觸發】daily_skill_report 開始執行")
    with open("C:/Users/blue_/claude-telegram-bot/schedule.log", "a", encoding="utf-8") as f:
        f.write(f"[{datetime.datetime.now().strftime('%Y-%m-%d %H:%M:%S')}] daily_skill_report 觸發\n")
    global _report_today
    import asyncio
    loop = asyncio.get_running_loop()

    today = datetime.date.today()
    today_str = str(today)

    if _report_today == today_str:
        return
    _report_today = today_str

    def count_tools():
        # 計算 TOOLS 列表中的工具數量
        return len(TOOLS)

    def count_today_commits():
        try:
            repo = str(Path(__file__).parent)
            result = subprocess.run(
                ["git", "-C", repo, "log", "--oneline", f"--since={today_str}", "--until={0}".format(
                    str(datetime.date.today() + datetime.timedelta(days=1))
                )],
                capture_output=True, text=True, encoding="utf-8"
            )
            commits = [l for l in result.stdout.strip().split("\n") if l.strip()]
            # 過濾出含 feat/fix/skill 的 commit（代表新增或修改技能）
            skill_commits = [c for c in commits if any(k in c.lower() for k in ["feat", "fix", "skill", "技能", "新增", "add"])]
            return len(skill_commits), commits
        except Exception:
            return 0, []

    total_tools = await loop.run_in_executor(None, count_tools)
    skill_count, all_commits = await loop.run_in_executor(None, count_today_commits)

    def generate_report():
        commit_list = "\n".join(all_commits[:10]) if all_commits else "今天沒有任何 commit"
        response = client.messages.create(
            model="claude-sonnet-4-6",
            max_tokens=200,
            system="你是小牛馬，說話嘴賤幽默。用繁體中文簡短回報今日技能安裝狀況，像跟老闆報告一樣，簡潔有力。",
            messages=[{"role": "user", "content": (
                f"今日技能報告：\n"
                f"- 目前總技能數：{total_tools} 個\n"
                f"- 今日技能相關更新：{skill_count} 筆\n"
                f"- 今日所有 commit：\n{commit_list}\n\n"
                f"用嘴賤風格簡短回報。"
            )}]
        )
        return response.content[0].text

    msg = await loop.run_in_executor(None, generate_report)
    await context.bot.send_message(chat_id=OWNER_ID, text=f"🔧 今日技能報告\n\n{msg}")


def collect_daily_report() -> str:
    """收集今天電腦的活動記錄，回傳給 Claude 整理成報告"""
    import subprocess, datetime
    lines = []
    today = datetime.date.today().strftime("%Y-%m-%d")

    # 1. 系統開關機 / 重開機事件
    try:
        ps = subprocess.run(
            ["powershell", "-Command",
             f"[Console]::OutputEncoding=[System.Text.Encoding]::UTF8;"
             f"Get-WinEvent -FilterHashtable @{{LogName='System';Id=@(1074,6005,6006,6008,41);StartTime='{today}'}} "
             f"-ErrorAction SilentlyContinue | Select-Object TimeCreated,Id,Message | ConvertTo-Csv -NoTypeInformation"],
            capture_output=True, text=True, encoding="utf-8"
        )
        if ps.stdout.strip():
            lines.append("【開關機/重開機事件】\n" + ps.stdout.strip()[:800])
    except Exception:
        pass

    # 2. 安裝或更新的程式
    try:
        ps = subprocess.run(
            ["powershell", "-Command",
             f"[Console]::OutputEncoding=[System.Text.Encoding]::UTF8;"
             f"Get-WinEvent -FilterHashtable @{{LogName='System';Id=@(19,20,43);StartTime='{today}'}} "
             f"-ErrorAction SilentlyContinue | Select-Object TimeCreated,Message | ConvertTo-Csv -NoTypeInformation"],
            capture_output=True, text=True, encoding="utf-8"
        )
        if ps.stdout.strip():
            lines.append("【程式安裝/更新】\n" + ps.stdout.strip()[:500])
    except Exception:
        pass

    # 3. 今天執行過的程式（Security 事件 4688）
    try:
        ps = subprocess.run(
            ["powershell", "-Command",
             f"[Console]::OutputEncoding=[System.Text.Encoding]::UTF8;"
             f"Get-WinEvent -FilterHashtable @{{LogName='Security';Id=4688;StartTime='{today}'}} "
             f"-ErrorAction SilentlyContinue -MaxEvents 50 | "
             f"Select-Object TimeCreated,@{{N='Process';E={{($_.Message -split '新增處理程序名稱:')[1] -split '\`n' | Select-Object -First 1}}}} | "
             f"ConvertTo-Csv -NoTypeInformation"],
            capture_output=True, text=True, encoding="utf-8"
        )
        if ps.stdout.strip():
            lines.append("【今日執行程式（前50筆）】\n" + ps.stdout.strip()[:800])
    except Exception:
        pass

    # 4. 今天詢問小牛馬的對話數據表（從 memory.db）
    try:
        import sqlite3
        db_path = Path(__file__).parent / "memory.db"
        if db_path.exists():
            conn = sqlite3.connect(str(db_path))
            cur = conn.cursor()

            # 今日所有訊息
            cur.execute(
                "SELECT chat_id, role, content, created_at FROM chat_history WHERE created_at >= ? ORDER BY created_at",
                (today,)
            )
            rows = cur.fetchall()
            conn.close()

            if rows:
                # 統計各維度
                total_user = sum(1 for r in rows if r[1] == "user")
                total_bot  = sum(1 for r in rows if r[1] == "assistant")

                # 每個 chat_id 的問題數
                from collections import Counter
                chat_counts = Counter(r[0] for r in rows if r[1] == "user")

                # 時段分布（每2小時一格）
                hour_counts = Counter()
                for r in rows:
                    if r[1] == "user":
                        try:
                            h = int(r[3][11:13])  # HH from "YYYY-MM-DD HH:MM:SS"
                            slot = f"{h:02d}:00~{h+2:02d}:00" if h % 2 == 0 else f"{h-1:02d}:00~{h+1:02d}:00"
                            hour_counts[f"{(h//2)*2:02d}:00"] += 1
                        except Exception:
                            pass

                # 用戶訊息清單（給 Claude 分析主題用）
                user_msgs = [r[2][:80] for r in rows if r[1] == "user"]

                table = (
                    f"【小牛馬今日對話數據表】\n"
                    f"{'─'*30}\n"
                    f"用戶提問總數：{total_user} 則\n"
                    f"小牛馬回覆數：{total_bot} 則\n"
                    f"活躍聊天室數：{len(chat_counts)} 個\n"
                    f"\n各聊天室提問數：\n"
                    + "\n".join(f"  chat_id {cid}：{cnt} 則" for cid, cnt in chat_counts.most_common())
                    + f"\n\n時段分布：\n"
                    + "\n".join(f"  {slot}：{'█' * cnt} {cnt}則" for slot, cnt in sorted(hour_counts.items()))
                    + f"\n\n今日問題內容（供分析主題）：\n"
                    + "\n".join(f"  - {m}" for m in user_msgs[:20])
                )
                lines.append(table)
            else:
                lines.append("【小牛馬今日對話數據表】\n今日無對話記錄")
    except Exception as e:
        lines.append(f"【小牛馬今日對話數據表】讀取失敗：{e}")

    # 5. 目前系統資源
    try:
        import psutil
        cpu = psutil.cpu_percent(interval=1)
        mem = psutil.virtual_memory()
        disk = psutil.disk_usage("C:\\")
        lines.append(
            f"【目前系統狀態】\n"
            f"CPU：{cpu}%　記憶體：{mem.percent}%（{mem.used//1024//1024//1024:.1f}GB / {mem.total//1024//1024//1024:.1f}GB）\n"
            f"磁碟C：使用 {disk.percent}%（剩餘 {disk.free//1024//1024//1024:.1f}GB）"
        )
    except Exception:
        pass

    return "\n\n".join(lines) if lines else "今日無特殊活動記錄"


async def daily_pc_report(context: ContextTypes.DEFAULT_TYPE):
    """每天晚上 10:00 台灣時間：匯報今天電腦做了什麼"""
    logging.info("【排程觸發】daily_pc_report 開始執行")
    with open("C:/Users/blue_/claude-telegram-bot/schedule.log", "a", encoding="utf-8") as f:
        f.write(f"[{datetime.datetime.now().strftime('%Y-%m-%d %H:%M:%S')}] daily_pc_report 觸發\n")
    import asyncio
    loop = asyncio.get_running_loop()

    raw_data = await loop.run_in_executor(None, collect_daily_report)

    def generate_report():
        response = client.messages.create(
            model="claude-sonnet-4-6",
            max_tokens=1000,
            system=(
                "你是小牛馬，于晏哥的貼身助理。現在是晚上10點，你要把今天這台電腦的活動整理成一份有條理的日報給于晏哥。"
                "風格要簡潔專業，用繁體中文條列式呈現，分成幾個區塊。如果沒什麼特別的就直說，不要廢話。"
                "結尾可以加一句嘴賤的話。"
            ),
            messages=[{"role": "user", "content": f"以下是今天的電腦活動原始資料，請整理成日報：\n\n{raw_data}"}]
        )
        return response.content[0].text

    report = await loop.run_in_executor(None, generate_report)
    await context.bot.send_message(chat_id=OWNER_ID, text=f"📋 今日電腦日報\n\n{report}")


if __name__ == "__main__":
    token = os.getenv("TELEGRAM_BOT_TOKEN")
    app = ApplicationBuilder().token(token).build()

    # 每天早上 11:00 台灣時間（UTC+8）= 03:00 UTC
    app.job_queue.run_daily(
        goodmorning,
        time=datetime.time(hour=3, minute=0, tzinfo=datetime.timezone.utc)
    )
    # 每天晚上 10:00 台灣時間（UTC+8）= 14:00 UTC
    app.job_queue.run_daily(
        daily_pc_report,
        time=datetime.time(hour=14, minute=0, tzinfo=datetime.timezone.utc)
    )
    # 每天晚上 10:30 台灣時間（UTC+8）= 14:30 UTC
    app.job_queue.run_daily(
        goodnight,
        time=datetime.time(hour=14, minute=30, tzinfo=datetime.timezone.utc)
    )
    # 每天晚上 9:00 台灣時間（UTC+8）= 13:00 UTC
    app.job_queue.run_daily(
        daily_skill_report,
        time=datetime.time(hour=13, minute=0, tzinfo=datetime.timezone.utc)
    )

    app.add_handler(CommandHandler("start", start))
    app.add_handler(CommandHandler("chatid", chatid))
    app.add_handler(CommandHandler("myid", myid))
    app.add_handler(CommandHandler("clear", clear))
    app.add_handler(CommandHandler("memories", memories))
    app.add_handler(MessageHandler(filters.TEXT & ~filters.COMMAND, handle_message))
    app.add_handler(MessageHandler(filters.PHOTO, handle_photo_message, block=True))
    app.add_handler(MessageHandler(filters.VOICE | filters.AUDIO, handle_voice_message))

    import asyncio as _asyncio

    async def _run():
        async with app:
            await app.start()
            await app.updater.start_polling(drop_pending_updates=True)
            # 永久等待直到程序被終止
            while True:
                await _asyncio.sleep(3600)

    try:
        _asyncio.run(_run())
    except (KeyboardInterrupt, SystemExit):
        pass
